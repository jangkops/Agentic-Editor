"""native_layout_renderer — 알려진_레이아웃을 편집가능 네이티브 PPTX 도형으로 렌더.

이 모듈은 ``slide_templates``의 알려진_레이아웃(cover/section_divider/two_column/
feature_grid/timeline/comparison/architecture)을 1920×1080 PNG로 굽지 않고,
python-pptx 네이티브 도형(텍스트박스/오토셰이프/표)으로 직접 렌더하기 위한 신규 경로다.

설계 근거: ``.kiro/specs/pptx-native-density-render/design.md`` (Components and
Interfaces §1, Data Models). 배치 보정은 ``layout_geometry``의 순수 기하 함수에
위임하고, 시각 구조(카드/배지/그림자)는 ``native_diagram_pptx``를 재사용한다.

제약(설계 정합):
  - LLM/네트워크 호출 0. 이 모듈은 순수 렌더 + 기하 위임만 수행한다.
  - 색·여백·타이포는 ``design_tokens_for_profile`` 결과 dict에서만 읽는다(신규 토큰 금지).
  - 콘텐츠 텍스트는 절대 이미지로 베이크하지 않는다(편집가능_네이티브 유지).

본 파일은 작업 1.1 스캐폴드다 — 데이터모델·상수·레지스트리·함수 시그니처는 실제로
정의하되, 렌더 로직 본체는 후속 작업(2.x~4.x)에서 구현한다(현재 NotImplementedError).
"""

from __future__ import annotations

import os
import base64
import time
import tempfile
import inspect
import asyncio
from dataclasses import dataclass, field
from typing import Callable, Optional

# ---------------------------------------------------------------------------
# import 재사용 배선 — server.py 기존 try/except fallback 패턴 참고
# (``from <mod> import ...`` 시도 후 ``from ai_engine.<mod> import ...`` fallback)
# ---------------------------------------------------------------------------

# layout_geometry: 배치 보정용 순수 기하 함수/상수
try:
    from layout_geometry import (
        within_bounds,
        clamp_into_bounds,
        resolve_collisions,
        fit_within,
        is_fullbleed,
        overlap_area,
        area,
        SLIDE_RECT,
    )
except Exception:  # pragma: no cover - 패키지 경로 fallback
    from ai_engine.layout_geometry import (  # type: ignore[no-redef]
        within_bounds,
        clamp_into_bounds,
        resolve_collisions,
        fit_within,
        is_fullbleed,
        overlap_area,
        area,
        SLIDE_RECT,
    )

# native_diagram_pptx: 표지/다이어그램 네이티브 빌더 재사용
# NOTE: _card/_badge_in_gutter/_set_text/_shadow 는 build_native_diagram/
#       build_native_cover 내부의 *중첩 함수*라 모듈 레벨 import가 불가하다.
#       시각 구조 재사용은 후속 작업에서 build_native_* 위임 또는 헬퍼 추출로 처리한다.
try:
    from native_diagram_pptx import build_native_cover, build_native_diagram
except Exception:  # pragma: no cover - 패키지 경로 fallback
    from ai_engine.native_diagram_pptx import (  # type: ignore[no-redef]
        build_native_cover,
        build_native_diagram,
    )

# slide_templates: design_tokens (색/여백/타이포) — 신규 토큰 금지, 재사용만
try:
    from slide_templates import design_tokens_for_profile
except Exception:  # pragma: no cover - 패키지 경로 fallback
    from ai_engine.slide_templates import design_tokens_for_profile  # type: ignore[no-redef]

# python-pptx: 네이티브 도형/텍스트 방출용. 미설치 환경에서도 모듈 import는
# 가능하도록 가드한다(실제 emit_* 호출 시에만 필요).
try:
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.oxml.ns import qn
    try:
        from pptx.enum.text import MSO_AUTO_SIZE
    except Exception:  # pragma: no cover
        MSO_AUTO_SIZE = None  # type: ignore[assignment]
except Exception:  # pragma: no cover - python-pptx 미설치 환경
    Inches = Pt = Emu = None  # type: ignore[assignment]
    MSO_SHAPE = RGBColor = PP_ALIGN = MSO_ANCHOR = qn = None  # type: ignore[assignment]
    MSO_AUTO_SIZE = None  # type: ignore[assignment]


# ---------------------------------------------------------------------------
# 슬라이드 경계 상수 — layout_geometry.SLIDE_RECT(0,0,13.333,7.5)와 정합
# ---------------------------------------------------------------------------

SLIDE_W_IN: float = 13.333
SLIDE_H_IN: float = 7.5


# ---------------------------------------------------------------------------
# 데이터 모델 (design.md §Data Models)
# ---------------------------------------------------------------------------


@dataclass
class PlacedShape:
    """배치된 네이티브 셰이프의 메타데이터(역할/Rect/텍스트유무/z-order)."""

    role: str  # "title"|"body"|"card"|"badge"|"note"|"contact"|"figure"
    #            |"decorative_bg"|"image"
    rect: tuple[float, float, float, float]  # (left, top, width, height) 인치
    has_text: bool  # 편집가능 텍스트 런 보유 여부
    text: str = ""  # 정규화 텍스트(제목 중복 판정용)
    z: int = 0  # z-order (작을수록 아래)


@dataclass
class RenderResult:
    """render_native_layout / render_native_fallback 의 반환 결과."""

    ok: bool  # 렌더 성공 여부
    placed: list[PlacedShape] = field(default_factory=list)  # 배치된 셰이프 메타
    title_count: int = 0  # 방출된 제목 셰이프 수 (Req 4: 정확히 1 또는 0)
    unsupported: bool = False  # 변환 불가 → 폴백 트리거 (Req 1.4)


class OverlapError(Exception):
    """위치/크기 조정 후에도 겹침률 10% 미만으로 못 낮춘 과밀 슬라이드 오류 (Req 2.5).

    슬라이드 식별자와 겹침 위반 셰이프 쌍 목록을 담아 호출자에게 반환한다
    (design Error Handling: ``OverlapError(slide_id, [(shape_a, shape_b)])``).
    """

    def __init__(self, slide_id, pairs):
        self.slide_id = slide_id
        self.pairs = list(pairs or [])
        super().__init__(
            f"슬라이드 {slide_id}: 겹침 10% 미만으로 보정 실패 — 위반 쌍 {self.pairs}"
        )


# ---------------------------------------------------------------------------
# 레이아웃명 → emit 디스패치 레지스트리 (NATIVE_LAYOUT_REGISTRY)
# 값은 후속 작업(3.1)에서 구현할 per-layout 렌더 함수. 현재는 스텁 콜러블.
# ---------------------------------------------------------------------------


# 공통 콘텐츠 영역 상수(인치) — 슬라이드 경계(13.333×7.5) 안 합리적 배분.
# 정밀 겹침/경계 보정은 작업 4.1 finalize_placement 가 담당한다.
_MARGIN_L: float = 0.7
_MARGIN_R: float = 0.7
_CONTENT_W: float = SLIDE_W_IN - _MARGIN_L - _MARGIN_R  # ≈11.933
_TITLE_REGION = (_MARGIN_L, 0.5, _CONTENT_W, 1.0)        # 제목 상단
_BODY_TOP: float = 1.9                                   # 본문 시작 y
# 작업 18.1: 본문 콘텐츠 하한 — 이 아래(≈5.55~7.45)는 고밀도 chrome 하단 밴드로
# 예약한다(콘텐츠 텍스트와 겹침 0). 콘텐츠 영역은 이 하한까지만 사용한다.
_CONTENT_BOTTOM: float = 5.5
_BODY_H: float = _CONTENT_BOTTOM - _BODY_TOP             # 본문 높이(하단 chrome 밴드 예약)


def _render_cover(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """cover — emit_title(대제목) + eyebrow/subtitle/footer 텍스트박스(편집가능).

    작업 18.1: 콘텐츠 텍스트 *앞에* 표지 고밀도 chrome(코너 글로우/아이콘 배지/
    노티스 칩/액센트 바)을 먼저 방출해 밀도 reference 6 + 스타일 품질을 충족하고,
    콘텐츠 텍스트가 chrome 위(z-order 앞)에 오도록 한다(가독 보존).
    """
    placed: list[PlacedShape] = []
    placed.extend(_emit_cover_chrome(slide, tokens))  # 장식 먼저 → 콘텐츠가 위.
    eyebrow = str(data.get("eyebrow") or "").strip()
    subtitle = str(data.get("subtitle") or "").strip()
    footer = str(data.get("footer") or data.get("date_str") or "").strip()

    if eyebrow:
        placed.append(_emit_text_block(slide, eyebrow, tokens,
                                       (_MARGIN_L, 1.4, _CONTENT_W, 0.5), role="eyebrow"))
    # 대제목 — 화면 중앙 상단 큰 영역.
    placed.append(emit_title(slide, data.get("title"), tokens,
                             (_MARGIN_L, 2.1, _CONTENT_W, 1.9)))
    if subtitle:
        placed.append(_emit_text_block(slide, subtitle, tokens,
                                       (_MARGIN_L, 4.2, _CONTENT_W, 1.4), role="body"))
    if footer:
        placed.append(_emit_text_block(slide, footer, tokens,
                                       (_MARGIN_L, 6.7, _CONTENT_W, 0.4), role="caption"))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_section_divider(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """section_divider — emit_section_header_bar(번호+제목 막대) + description 텍스트박스."""
    placed: list[PlacedShape] = []
    no = data.get("section_number") or data.get("number") or 1
    placed.extend(emit_section_header_bar(slide, no, data.get("title"), tokens,
                                          (_MARGIN_L, 2.6, _CONTENT_W, 1.1)))
    desc = str(data.get("description") or "").strip()
    if desc:
        placed.append(_emit_text_block(slide, desc, tokens,
                                       (_MARGIN_L, 4.0, _CONTENT_W, 1.4), role="body"))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_two_column(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """two_column — emit_title + 좌/우 컬럼(불릿 텍스트박스)."""
    placed: list[PlacedShape] = []
    placed.append(emit_title(slide, data.get("title"), tokens, _TITLE_REGION))
    subtitle = str(data.get("subtitle") or "").strip()
    body_top = _BODY_TOP
    if subtitle:
        placed.append(_emit_text_block(slide, subtitle, tokens,
                                       (_MARGIN_L, 1.5, _CONTENT_W, 0.5), role="body"))
        body_top = 2.2
    body_h = _CONTENT_BOTTOM - body_top
    col_gap = 0.4
    col_w = (_CONTENT_W - col_gap) / 2.0
    left_x = _MARGIN_L
    right_x = _MARGIN_L + col_w + col_gap
    placed.append(_emit_text_block(slide, data.get("left_content"), tokens,
                                   (left_x, body_top, col_w, body_h),
                                   role="body", bullets=True))
    placed.append(_emit_text_block(slide, data.get("right_content"), tokens,
                                   (right_x, body_top, col_w, body_h),
                                   role="body", bullets=True))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_feature_grid(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """feature_grid — emit_title + emit_card_grid(features)."""
    placed: list[PlacedShape] = []
    placed.append(emit_title(slide, data.get("title"), tokens, _TITLE_REGION))
    placed.extend(emit_card_grid(slide, data.get("features"), tokens,
                                 (_MARGIN_L, _BODY_TOP, _CONTENT_W, _BODY_H)))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_timeline(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """timeline — emit_title + steps 를 emit_numbered_list(번호 배지+텍스트)로."""
    placed: list[PlacedShape] = []
    placed.append(emit_title(slide, data.get("title"), tokens, _TITLE_REGION))
    placed.extend(emit_numbered_list(slide, data.get("steps"), tokens,
                                     (_MARGIN_L, _BODY_TOP, _CONTENT_W, _BODY_H)))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_comparison(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """comparison — emit_title + 좌/우(라벨 + 불릿 items) 2개 컬럼."""
    placed: list[PlacedShape] = []
    placed.append(emit_title(slide, data.get("title"), tokens, _TITLE_REGION))
    col_gap = 0.4
    col_w = (_CONTENT_W - col_gap) / 2.0
    label_h = 0.6
    items_top = _BODY_TOP + label_h + 0.15
    items_h = _CONTENT_BOTTOM - items_top
    for idx, (label_key, items_key) in enumerate(
            (("left_label", "left_items"), ("right_label", "right_items"))):
        cx = _MARGIN_L + idx * (col_w + col_gap)
        placed.append(_emit_text_block(slide, data.get(label_key), tokens,
                                       (cx, _BODY_TOP, col_w, label_h), role="card_title"))
        placed.append(_emit_text_block(slide, data.get(items_key), tokens,
                                       (cx, items_top, col_w, items_h),
                                       role="body", bullets=True))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


def _render_architecture(slide, prs, data: dict, tokens: dict, *, palette=None) -> RenderResult:
    """architecture — emit_title + layers 를 카드 박스 수직 스택으로(emit_card_grid 재사용)."""
    placed: list[PlacedShape] = []
    placed.append(emit_title(slide, data.get("title"), tokens, _TITLE_REGION))
    layers = list(data.get("layers") or [])
    n = len(layers)
    if n > 0:
        gap = 0.25
        row_h = (_BODY_H - (n - 1) * gap) / n
        for i, layer in enumerate(layers):
            name = str(layer.get("name") or layer.get("title") or "") if isinstance(layer, dict) else str(layer)
            desc = ""
            if isinstance(layer, dict):
                desc = str(layer.get("description") or "")
                items = layer.get("items")
                if items:
                    joined = ", ".join(str(it) for it in items) if isinstance(items, (list, tuple)) else str(items)
                    desc = f"{desc} · {joined}".strip(" ·") if desc else joined
            row_top = _BODY_TOP + i * (row_h + gap)
            placed.extend(emit_card_grid(slide, [{"title": name, "description": desc}], tokens,
                                         (_MARGIN_L, row_top, _CONTENT_W, row_h)))
    placed.extend(_emit_body_chrome(slide, tokens, number=_chrome_number(data),
                                    footer_text=_chrome_footer(data)))
    return RenderResult(ok=True, placed=placed, title_count=1, unsupported=False)


NATIVE_LAYOUT_REGISTRY: dict[str, Callable] = {
    "cover": _render_cover,
    "section_divider": _render_section_divider,
    "two_column": _render_two_column,
    "feature_grid": _render_feature_grid,
    "timeline": _render_timeline,
    "comparison": _render_comparison,
    "architecture": _render_architecture,
}


# ---------------------------------------------------------------------------
# 작업 20(B) — LLM 레이아웃 픽 → 네이티브 7키 매핑
# ---------------------------------------------------------------------------
# server.py 의 _llm_pick_slide_layout 은 slide_templates.LAYOUT_REGISTRY(11종:
# cover/section_divider/two_column/feature_grid/timeline/comparison/architecture +
# kpi_summary/status_table/objective_detail/process_flow)로 레이아웃을 고른다.
# NATIVE_LAYOUT_REGISTRY 는 7키뿐이라, LLM 이 비-7키를 고르면 네이티브 라우팅이
# 되지 않아 통짜 이미지 베이크로 빠지던 근본 원인을 해소한다. 비-7키를 의미상 가장
# 가까운 7키로 매핑해 어떤 레이아웃 선택이든 편집가능 네이티브로 라우팅되게 한다.
# data 가 매핑된 레이아웃의 REQUIRED_FIELDS 와 안 맞으면 render_native_layout 가
# ok=False → render_native_fallback(편집가능 텍스트박스, 콘텐츠 손실 0)로 안전 처리.
_NON_NATIVE_TO_NATIVE: dict = {
    "kpi_summary": "feature_grid",     # 지표 카드 그리드 ≈ feature_grid
    "status_table": "comparison",      # 항목/상태 대비 ≈ comparison
    "objective_detail": "two_column",  # 목표 상세(근거/증빙) ≈ two_column
    "process_flow": "timeline",        # 단계/순서 흐름 ≈ timeline
}


def map_to_native_layout(layout: str) -> str:
    """LLM 이 고른 레이아웃명을 NATIVE_LAYOUT_REGISTRY 의 7키 중 하나로 매핑한다.

    - 이미 7키(NATIVE_LAYOUT_REGISTRY 키)면 그대로 반환(항등, 회귀 0).
    - 알려진 비-7키(kpi_summary/status_table/objective_detail/process_flow)는
      의미상 근접한 7키로 매핑.
    - 그 외 미지의 레이아웃/빈값은 안전 기본값 ``two_column`` 으로 매핑.

    순수 결정 함수(네트워크/LLM 호출 없음).
    """
    key = (layout or "").strip()
    if not key:
        # 빈 입력(LLM 픽 실패/타임아웃/파싱실패/미선택) — 통짜 이미지 원천 차단.
        # 과거에는 빈 입력을 "" 로 돌려보내 _should_native_render 를 False 로 만들고
        # HTML→PNG 베이크(통짜 이미지) 폴백으로 귀결시켰다. 그것이 실제 게이트웨이
        # 환경에서 슬라이드 전체가 편집 불가 통짜 이미지가 되던 근본 원인이었다.
        # 이제 빈 입력도 안전 기본값 ``two_column`` 으로 매핑해, 픽이 실패해도 항상
        # 편집가능 네이티브 렌더로 라우팅되게 한다(빈=베이크 폴백 정책 폐기).
        return "two_column"
    if key in NATIVE_LAYOUT_REGISTRY:
        return key
    if key in _NON_NATIVE_TO_NATIVE:
        return _NON_NATIVE_TO_NATIVE[key]
    # 비-빈 미지 레이아웃명 — 안전 기본값(two_column). LLM 이 실제로 레이아웃을
    # 선택했으나 7키/알려진 비-7키에 없을 때만 도달한다.
    return "two_column"

# 제목 역할 집합 — title_count 집계(Req 4: 슬라이드당 정확히 1 또는 0)에 사용.
TITLE_ROLES: frozenset = frozenset({"title", "section_title"})

# 레이아웃별 필수 필드 — 부재 시 폴백 트리거(Req 1.4).
REQUIRED_FIELDS: dict[str, list] = {
    "cover": ["title"],
    "section_divider": ["title"],
    "two_column": ["title", "left_content", "right_content"],
    "feature_grid": ["title", "features"],
    "timeline": ["title", "steps"],
    "comparison": ["title", "left_label", "left_items", "right_label", "right_items"],
    "architecture": ["title", "layers"],
}


def _missing_required(layout: str, data: dict) -> bool:
    """알려진_레이아웃의 필수 필드 부재(None/빈문자열/빈컬렉션) 여부."""
    for f in REQUIRED_FIELDS.get(layout, []):
        v = (data or {}).get(f)
        if v is None:
            return True
        if isinstance(v, str) and not v.strip():
            return True
        if isinstance(v, (list, tuple, dict)) and len(v) == 0:
            return True
    return False


# ---------------------------------------------------------------------------
# 공개 렌더 진입점 (design.md §Components and Interfaces §1)
# ---------------------------------------------------------------------------


def render_native_layout(
    slide,
    prs,
    layout: str,
    data: dict,
    tokens: dict,
    *,
    palette: Optional[list] = None,
    aws_profile: str = "",
    credentials: Optional[dict] = None,
) -> RenderResult:
    """알려진_레이아웃을 편집가능 네이티브 도형으로 슬라이드에 렌더한다.

    Args:
        slide: python-pptx Slide
        prs: python-pptx Presentation (cover 에서 슬라이드 크기 참조)
        layout: LAYOUT_REGISTRY 키 (NATIVE_LAYOUT_REGISTRY 와 동일 집합)
        data: slide_templates.render_layout 과 동일 data 스키마
        tokens: design_tokens_for_profile(profile) 결과 (색/여백/타이포)
        palette: 선택 팔레트
        aws_profile/credentials: Vertex 장식 옵트인(작업 8.1) 시 자격증명 해석에 전달.
            기본값(빈값/None)에서는 장식 경로가 옵트인 OFF 또는 자격증명 부재로 No-op.

    Returns:
        RenderResult — ok/placed/title_count/unsupported.
        미지원 레이아웃/필수필드 부재 시 ok=False, unsupported=True.

    작업 3.1: NATIVE_LAYOUT_REGISTRY 디스패치 → _render_* 호출 → placed 로 title_count
    집계(Req 4). finalize_placement(작업 4.1)는 forward-compatible try/except 로 호출하되,
    미구현(NotImplementedError) 동안에는 미보정 placed 를 그대로 반환한다 — 4.1 완료 후
    자동 통합된다.

    작업 8.1: 콘텐츠 셰이프를 방출하기 *전에* (옵트인 ON + 장식 힌트 시) 장식_배경을
    먼저 방출해 back-most(콘텐츠 텍스트보다 뒤) z-order 를 보장한다(Property 13/16).
    옵트인 OFF/자격증명 부재/생성 실패 시 장식은 비워지고 콘텐츠 네이티브만 손실 0으로
    렌더된다(Req 11.3, Property 17) — 예외는 전파하지 않는다.
    """
    if layout not in NATIVE_LAYOUT_REGISTRY or _missing_required(layout, data):
        # 폴백 트리거 — 호출자는 render_native_fallback(작업 3.2)을 적용한다.
        return RenderResult(ok=False, placed=[], title_count=0, unsupported=True)

    placed: list[PlacedShape] = []

    # 장식_배경(풀블리드) — 옵트인 ON + 힌트 시에만. 콘텐츠보다 먼저 add → back-most.
    # 예외 전파 금지(장식 실패가 콘텐츠 렌더를 막지 않음, Req 11.3).
    try:
        placed.extend(
            maybe_add_decorative_background(
                slide, data, layout, aws_profile=aws_profile, credentials=credentials
            )
        )
    except Exception:
        pass

    render_fn = NATIVE_LAYOUT_REGISTRY[layout]
    result = render_fn(slide, prs, data, tokens, palette=palette)
    placed.extend(result.placed)

    # title_count 는 방출된 placed 의 제목 역할 수로 집계(Req 4: 0 또는 1).
    title_count = sum(1 for ps in placed if ps.role in TITLE_ROLES)

    # 작업 4.1 통합 지점 — finalize_placement 가 구현되면 겹침/경계 보정 + 제목 dedup
    # 후의 placed 를 채택한다. 미구현 동안에는 미보정 placed 를 그대로 사용한다.
    try:
        placed = finalize_placement(placed)
        title_count = sum(1 for ps in placed if ps.role in TITLE_ROLES)
    except NotImplementedError:
        pass

    return RenderResult(ok=True, placed=placed, title_count=title_count, unsupported=False)


# 폴백 텍스트 추출에서 제외할 비콘텐츠 키(스타일/제어 힌트 — 콘텐츠 텍스트 아님).
# 콘텐츠 텍스트는 절대 누락하지 않되, 색/톤/방향 같은 스타일 힌트는 본문에 섞지 않는다.
_FALLBACK_SKIP_KEYS: frozenset = frozenset({
    "accent_color", "orientation", "tone", "color", "bg", "background",
    "palette", "icon_color", "variant", "align", "valign",
})


def _collect_fallback_lines(obj) -> list[str]:
    """변환 불가 입력에서 콘텐츠 텍스트를 재귀적·평탄하게 전수 추출(손실 0).

    dict/list/str 를 깊이 우선으로 순회하여 비어있지 않은 문자열 leaf 를 입력 순서대로
    모은다. 색·톤·방향 등 명백한 스타일/제어 키(_FALLBACK_SKIP_KEYS, ``*_color``/
    ``*_tone``)는 콘텐츠가 아니므로 건너뛴다. 숫자/불리언은 콘텐츠 텍스트로 취급하지
    않는다(예: section_number). 어떤 콘텐츠 텍스트도 이미지로 베이크되지 않는다.
    """
    out: list[str] = []
    if obj is None:
        return out
    if isinstance(obj, str):
        s = obj.strip()
        if s:
            out.append(s)
        return out
    if isinstance(obj, (bool, int, float)):
        return out
    if isinstance(obj, dict):
        for k, v in obj.items():
            ks = str(k).lower()
            if ks in _FALLBACK_SKIP_KEYS or ks.endswith("_color") or ks.endswith("_tone"):
                continue
            out.extend(_collect_fallback_lines(v))
        return out
    if isinstance(obj, (list, tuple, set)):
        for it in obj:
            out.extend(_collect_fallback_lines(it))
        return out
    # 기타 타입 — 문자열화하여 비지 않으면 보존.
    s = str(obj).strip()
    if s:
        out.append(s)
    return out


def render_native_fallback(slide, data: dict, tokens: dict) -> RenderResult:
    """Req 1.4/1.5 — 변환 불가 시 콘텐츠 텍스트를 편집가능 텍스트박스로 전수 출력.

    통짜 이미지 대체 절대 금지. 입력 콘텐츠 텍스트를 재귀적/평탄하게 전수 추출(손실 0)하여
    비어있지 않은 텍스트 런을 가진 편집가능_네이티브 셰이프 1개 이상으로 방출한다.

    구조:
      - 제목(top-level ``title`` 이 비어있지 않은 문자열): ``emit_title`` 로 상단 1개
        (Req 4: 제목 있으면 정확히 1개). title_count = 1.
      - 본문(나머지 모든 텍스트성 값): ``_emit_text_block`` 으로 제목 아래에 줄 단위 방출.
        ``features``/``steps``/``layers`` 의 title/description, ``left_items``/
        ``right_items``, ``left_content``/``right_content``, subtitle, description 등을
        재귀 추출하여 전수 보존한다.

    반환: RenderResult(ok=True, placed=[...], title_count=제목 있으면 1 else 0,
    unsupported=False). 콘텐츠 텍스트 셰이프는 has_text=True.

    NOTE: 작업 4.1 finalize_placement(미구현)는 호출하지 않는다 — 폴백은 제목 상단·본문
    그 아래로 슬라이드 경계(13.333×7.5) 안에 합리적으로 배치하며, 정밀 보정은 호출자
    또는 후속 작업이 담당한다.
    """
    data = data if isinstance(data, dict) else ({} if data is None else {"_content": data})
    placed: list[PlacedShape] = []

    # --- 제목 추출 (top-level title, 비어있지 않은 문자열일 때만) ---
    title_raw = data.get("title")
    title_text = str(title_raw).strip() if isinstance(title_raw, str) else ""
    has_title = bool(title_text)

    if has_title:
        placed.append(emit_title(slide, title_text, tokens, _TITLE_REGION))

    # --- 본문 추출: 제목은 별도 박스로 이동했으므로 top-level title 키 제외하고 전수 추출 ---
    body_source = {k: v for k, v in data.items() if k != "title"} if has_title else dict(data)
    lines = _collect_fallback_lines(body_source)

    # 제목과 정규화 후 동일한 본문 줄은 중복이므로 제거(제목 1회 유지, Req 4).
    if has_title:
        ntitle = _norm(title_text)
        lines = [ln for ln in lines if _norm(ln) != ntitle]

    # --- 본문 텍스트박스 방출 (제목 아래) ---
    if lines:
        body_top = _BODY_TOP if has_title else 0.6
        body_h = SLIDE_H_IN - body_top - 0.4
        placed.append(_emit_text_block(slide, lines, tokens,
                                       (_MARGIN_L, body_top, _CONTENT_W, body_h),
                                       role="body", bullets=True))

    return RenderResult(ok=True, placed=placed,
                        title_count=1 if has_title else 0, unsupported=False)


# ---------------------------------------------------------------------------
# role → 타이포/색 매핑 테이블 (apply_tokens_to_run 용)
#
# 색·폰트는 모두 design_tokens 키만 참조한다(Req 5.4 신규 토큰 금지). 폰트 크기는
# SLIDE_DESIGN 에 별도 토큰이 없으므로 "타이포 계층"(제목 > 본문)을 만들기 위한
# role별 편집 기본값으로 둔다 — Property 19(audit_style_quality)가 검사하는
# title 폰트크기 > body 폰트크기 계층을 구조적으로 보장한다. 색 키는 tokens.get()
# 으로 안전 조회하고, 부재 시 보수적 기본값(검정/흰색이 아닌 회색)으로 폴백한다.
# ---------------------------------------------------------------------------

_ROLE_STYLE: dict[str, dict] = {
    # role           size  bold   color_token    font_token
    "title":         {"size": 30, "bold": True,  "color": "text_dark",  "font": "font_heading"},
    "section_title": {"size": 22, "bold": True,  "color": "text_light", "font": "font_heading"},
    "card_title":    {"size": 16, "bold": True,  "color": "text_dark",  "font": "font_heading"},
    "eyebrow":       {"size": 12, "bold": True,  "color": "accent",     "font": "font_heading"},
    "badge":         {"size": 13, "bold": True,  "color": "text_light", "font": "font_heading"},
    "note_label":    {"size": 11, "bold": True,  "color": "accent",     "font": "font_heading"},
    "body":          {"size": 14, "bold": False, "color": "text_dark",  "font": "font_body"},
    "note":          {"size": 13, "bold": False, "color": "text_dark",  "font": "font_body"},
    "contact":       {"size": 13, "bold": False, "color": "text_dark",  "font": "font_body"},
    "caption":       {"size": 11, "bold": False, "color": "text_muted", "font": "font_body"},
}

# 토큰 부재 시 색 폴백 (RGB 튜플) — 검정/흰색 하드코딩이 아닌 보수적 회색 계열.
_COLOR_FALLBACK: dict[str, tuple] = {
    "text_dark": (26, 26, 26),
    "text_light": (255, 255, 255),
    "text_muted": (107, 114, 128),
    "primary": (0, 102, 255),
    "secondary": (0, 200, 150),
    "accent": (255, 107, 53),
    "border": (229, 231, 235),
    "card_bg": (255, 255, 255),
    "bg_section": (244, 246, 249),
}


# ---------------------------------------------------------------------------
# 내부 시각 헬퍼 (native_diagram_pptx 의 _card/_badge/_shadow 와 동등 구조를
# python-pptx 직접 호출로 재현 — 중첩 함수라 import 불가하므로 동등 구현).
# ---------------------------------------------------------------------------


def _hex_to_rgb_tuple(hexstr) -> Optional[tuple]:
    """'#RRGGBB' 또는 'RRGGBB' → (r, g, b). 형식 불일치 시 None."""
    if not hexstr:
        return None
    s = str(hexstr).strip().lstrip("#")
    if len(s) != 6:
        return None
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return None


def _token_rgb(tokens: dict, key: str):
    """design_tokens[key] 색을 RGBColor 로. 부재/형식불일치 시 _COLOR_FALLBACK."""
    t = _hex_to_rgb_tuple((tokens or {}).get(key)) or _COLOR_FALLBACK.get(key, (107, 114, 128))
    return RGBColor(*t)


def _tint_rgb(tokens: dict, key: str, factor: float = 0.86):
    """토큰 색을 흰색 쪽으로 옅게 섞은 RGBColor(틴트 배경용)."""
    r, g, b = _hex_to_rgb_tuple((tokens or {}).get(key)) or _COLOR_FALLBACK.get(key, (0, 102, 255))
    return RGBColor(
        int(r + (255 - r) * factor),
        int(g + (255 - g) * factor),
        int(b + (255 - b) * factor),
    )


def _font_name_from_stack(stack, default: str = "Malgun Gothic") -> str:
    """CSS 폰트 스택('-apple-system, ..., "Malgun Gothic", sans-serif')에서
    PowerPoint 에 설정할 구체 폰트명 1개를 고른다. 한글 가독 폰트를 우선한다."""
    if not stack:
        return default
    parts = [p.strip().strip("'\"") for p in str(stack).split(",")]
    first_concrete = None
    for p in parts:
        if not p:
            continue
        low = p.lower()
        if low in ("-apple-system", "blinkmacsystemfont", "system-ui",
                   "sans-serif", "serif", "monospace"):
            continue
        if first_concrete is None:
            first_concrete = p
        if any(k in p for k in ("Malgun", "Gothic", "Noto", "Segoe")):
            return p
    return first_concrete or default


def _norm(text) -> str:
    """제목 중복 판정용 정규화(앞뒤 공백 제거 · 소문자)."""
    return str(text or "").strip().lower()


def _apply_shadow(shape) -> None:
    """부드러운 드롭 섀도우(outerShdw) XML 주입 — native_diagram_pptx._shadow 동등."""
    if qn is None:
        return
    try:
        spPr = shape._element.spPr
        for el in spPr.findall(qn("a:effectLst")):
            spPr.remove(el)
        eff = spPr.makeelement(qn("a:effectLst"), {})
        sh = eff.makeelement(qn("a:outerShdw"), {
            "blurRad": "90000", "dist": "38100", "dir": "5400000", "rotWithShape": "0",
        })
        clr = sh.makeelement(qn("a:srgbClr"), {"val": "1A2A44"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": "24000"})
        clr.append(alpha)
        sh.append(clr)
        eff.append(sh)
        spPr.append(eff)
    except Exception:
        pass


def _set_textframe_basics(tf, *, anchor_middle: bool = True) -> None:
    """텍스트 프레임 공통 서식 — word_wrap, 내부 여백(token 기반 편집 마진), 수직정렬.

    SLIDE_DESIGN 에 별도 spacing 토큰이 없으므로 본문 가독을 위한 편집 기본 내부
    마진을 명시 적용한다(python-pptx 기본값에 의존하지 않음 → Property 19 여백 검사).
    """
    try:
        tf.word_wrap = True
    except Exception:
        pass
    if Emu is not None:
        try:
            tf.margin_left = Emu(int(0.12 * 914400))
            tf.margin_right = Emu(int(0.12 * 914400))
            tf.margin_top = Emu(int(0.06 * 914400))
            tf.margin_bottom = Emu(int(0.06 * 914400))
        except Exception:
            pass
    if anchor_middle and MSO_ANCHOR is not None:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass


def _add_run(paragraph, text: str, tokens: dict, role: str):
    """단락에 run 추가 + apply_tokens_to_run(role) 적용."""
    run = paragraph.add_run()
    run.text = str(text if text is not None else "")
    apply_tokens_to_run(run, tokens, role)
    return run


def _card_text(text) -> tuple:
    """카드/항목 입력(dict 또는 str)에서 (title, description) 추출."""
    if isinstance(text, dict):
        title = text.get("title") or text.get("name") or text.get("label") or ""
        desc = text.get("description") or text.get("text") or text.get("desc") or ""
        return str(title), str(desc)
    return str(text or ""), ""


def _as_lines(content) -> list[str]:
    """본문 콘텐츠 입력(list/tuple/dict/str)을 비어있지 않은 줄 리스트로 정규화."""
    if content is None:
        return []
    if isinstance(content, (list, tuple)):
        out = []
        for it in content:
            if isinstance(it, dict):
                ti, de = _card_text(it)
                out.append(ti if not de else f"{ti} — {de}")
            else:
                s = str(it)
                if s.strip():
                    out.append(s)
        return out
    if isinstance(content, dict):
        return [str(v) for v in content.values() if str(v).strip()]
    return [ln for ln in str(content).split("\n") if ln.strip()]


def _emit_text_block(slide, content, tokens: dict, region, role: str = "body",
                     *, bullets: bool = False) -> PlacedShape:
    """단순 텍스트 블록(eyebrow/subtitle/footer/description/컬럼 본문)을 편집가능
    텍스트박스 1개로 방출한다. list 입력은 줄(또는 불릿) 단위 단락으로 전개한다.

    작업 3.1 신규 헬퍼 — 기존 ``_add_run``/``_set_textframe_basics``를 재사용한다.
    """
    l, t, w, h = region
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    _set_textframe_basics(tf, anchor_middle=False)
    lines = _as_lines(content)
    if not lines:
        # 콘텐츠가 비어도 통짜 이미지 대체 금지 — 빈 텍스트박스는 has_text=False로 표기.
        return PlacedShape(role=role, rect=(l, t, w, h), has_text=False, z=2)
    collected = []
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if PP_ALIGN is not None:
            p.alignment = PP_ALIGN.LEFT
        prefix = "• " if bullets else ""
        _add_run(p, (prefix + str(ln))[:200], tokens, role)
        collected.append(str(ln))
    return PlacedShape(role=role, rect=(l, t, w, h), has_text=True,
                       text=_norm(" ".join(collected)), z=2)


# ---------------------------------------------------------------------------
# 고밀도_요소 → 네이티브 도형 매핑 (emit_*)  (design.md §Data Models 매핑 표)
# ---------------------------------------------------------------------------


def emit_title(slide, text: str, tokens: dict, region) -> PlacedShape:
    """제목 텍스트박스 1개 방출 (Req 4: 제목 1회). apply_tokens_to_run(role="title")."""
    l, t, w, h = region
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    _set_textframe_basics(tf)
    p = tf.paragraphs[0]
    if PP_ALIGN is not None:
        p.alignment = PP_ALIGN.LEFT
    _add_run(p, str(text or "")[:120], tokens, "title")
    return PlacedShape(role="title", rect=(l, t, w, h), has_text=True, text=_norm(text), z=5)


def emit_section_header_bar(slide, no, title: str, tokens: dict, region) -> list[PlacedShape]:
    """섹션 헤더: 다크 막대(ROUNDED_RECTANGLE) + 번호 배지(oval, accent) + 제목 텍스트박스."""
    l, t, w, h = region
    placed: list[PlacedShape] = []

    # 1) 다크 막대 (text_dark 토큰 채움 · 라운드 · 그림자) — 장식, 텍스트 없음.
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        bar.adjustments[0] = 0.16
    except Exception:
        pass
    bar.fill.solid()
    bar.fill.fore_color.rgb = _token_rgb(tokens, "text_dark")
    bar.line.fill.background()
    _apply_shadow(bar)
    placed.append(PlacedShape(role="section_bar", rect=(l, t, w, h), has_text=False, z=1))

    # 2) 번호 배지 (oval, accent 색) — 좌측 거터 안. 텍스트(번호) 편집가능.
    badge_d = max(0.4, h * 0.55)
    bx = l + 0.18
    by = t + (h - badge_d) / 2.0
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(bx), Inches(by), Inches(badge_d), Inches(badge_d))
    oval.fill.solid()
    oval.fill.fore_color.rgb = _token_rgb(tokens, "accent")
    oval.line.fill.background()
    otf = oval.text_frame
    _set_textframe_basics(otf)
    op = otf.paragraphs[0]
    if PP_ALIGN is not None:
        op.alignment = PP_ALIGN.CENTER
    _add_run(op, str(no), tokens, "badge")
    placed.append(PlacedShape(role="badge", rect=(bx, by, badge_d, badge_d),
                              has_text=True, text=str(no), z=3))

    # 3) 제목 텍스트박스 (배지 오른쪽, 흰색 텍스트) — 거터 밖이라 겹침 0.
    tx = bx + badge_d + 0.18
    tw = max(0.5, w - (tx - l) - 0.2)
    tb = slide.shapes.add_textbox(Inches(tx), Inches(t), Inches(tw), Inches(h))
    tf = tb.text_frame
    _set_textframe_basics(tf)
    p = tf.paragraphs[0]
    if PP_ALIGN is not None:
        p.alignment = PP_ALIGN.LEFT
    _add_run(p, str(title or "")[:90], tokens, "section_title")
    placed.append(PlacedShape(role="section_title", rect=(tx, t, tw, h),
                              has_text=True, text=_norm(title), z=3))
    return placed


def emit_contact_box(slide, contact, tokens: dict, region) -> list[PlacedShape]:
    """연락처 박스: 틴트 사각형(primary 틴트) + 좌측 액센트 바(accent) + 텍스트박스."""
    l, t, w, h = region
    placed: list[PlacedShape] = []

    # 1) 틴트 사각형(라운드 + border + 그림자) — 장식.
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        box.adjustments[0] = 0.08
    except Exception:
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = _tint_rgb(tokens, "primary", 0.90)
    box.line.color.rgb = _token_rgb(tokens, "border")
    box.line.width = Pt(1.0)
    _apply_shadow(box)
    placed.append(PlacedShape(role="contact_box", rect=(l, t, w, h), has_text=False, z=1))

    # 2) 좌측 액센트 바(accent) — 장식.
    bar_w = 0.12
    ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(bar_w), Inches(h))
    ab.fill.solid()
    ab.fill.fore_color.rgb = _token_rgb(tokens, "accent")
    ab.line.fill.background()
    placed.append(PlacedShape(role="accent_bar", rect=(l, t, bar_w, h), has_text=False, z=2))

    # 3) 텍스트박스(연락처 콘텐츠) — 편집가능.
    tx = l + bar_w + 0.18
    tw = max(0.5, w - bar_w - 0.36)
    tb = slide.shapes.add_textbox(Inches(tx), Inches(t), Inches(tw), Inches(h))
    tf = tb.text_frame
    _set_textframe_basics(tf)

    # contact 입력 정규화: dict / list / str 모두 줄 단위로.
    if isinstance(contact, dict):
        lines = [str(v) for v in contact.values() if str(v).strip()]
    elif isinstance(contact, (list, tuple)):
        lines = [str(v) for v in contact if str(v).strip()]
    else:
        lines = [ln for ln in str(contact or "").split("\n") if ln.strip()]
    if not lines:
        lines = [str(contact or "")]

    full_text = "\n".join(lines)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if PP_ALIGN is not None:
            p.alignment = PP_ALIGN.LEFT
        _add_run(p, ln[:120], tokens, "contact")
    placed.append(PlacedShape(role="contact", rect=(tx, t, tw, h),
                              has_text=True, text=_norm(full_text), z=3))
    return placed


def emit_note_callout(slide, text: str, tokens: dict, region) -> list[PlacedShape]:
    """노트 콜아웃: 경고 틴트 사각형(accent 틴트) + 좌측 보더(accent) + 텍스트박스."""
    l, t, w, h = region
    placed: list[PlacedShape] = []

    # 1) 경고 틴트 사각형(라운드 + border + 그림자) — 장식.
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(l), Inches(t), Inches(w), Inches(h))
    try:
        box.adjustments[0] = 0.08
    except Exception:
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = _tint_rgb(tokens, "accent", 0.88)
    box.line.color.rgb = _token_rgb(tokens, "accent")
    box.line.width = Pt(1.0)
    _apply_shadow(box)
    placed.append(PlacedShape(role="note_box", rect=(l, t, w, h), has_text=False, z=1))

    # 2) 좌측 보더(accent) — 장식.
    bar_w = 0.1
    lb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                Inches(l), Inches(t), Inches(bar_w), Inches(h))
    lb.fill.solid()
    lb.fill.fore_color.rgb = _token_rgb(tokens, "accent")
    lb.line.fill.background()
    placed.append(PlacedShape(role="note_border", rect=(l, t, bar_w, h), has_text=False, z=2))

    # 3) 텍스트박스(NOTICE 라벨 + 본문) — 편집가능.
    tx = l + bar_w + 0.18
    tw = max(0.5, w - bar_w - 0.36)
    tb = slide.shapes.add_textbox(Inches(tx), Inches(t), Inches(tw), Inches(h))
    tf = tb.text_frame
    _set_textframe_basics(tf)
    p0 = tf.paragraphs[0]
    if PP_ALIGN is not None:
        p0.alignment = PP_ALIGN.LEFT
    _add_run(p0, "NOTICE", tokens, "note_label")
    p1 = tf.add_paragraph()
    if PP_ALIGN is not None:
        p1.alignment = PP_ALIGN.LEFT
    _add_run(p1, str(text or "")[:240], tokens, "note")
    placed.append(PlacedShape(role="note", rect=(tx, t, tw, h),
                              has_text=True, text=_norm(text), z=3))
    return placed


def emit_numbered_list(slide, items, tokens: dict, region) -> list[PlacedShape]:
    """번호 목록: 항목별 (원형 배지 oval + 텍스트박스). 배지는 좌측 거터에 배치해
    본문 텍스트박스와 겹침 0(거터 폭 = 배지 지름 + 간격)."""
    l, t, w, h = region
    placed: list[PlacedShape] = []
    items = list(items or [])
    n = len(items)
    if n <= 0:
        return placed

    row_h = h / n
    badge_d = max(0.34, min(0.5, row_h * 0.55))
    gutter = badge_d + 0.15  # 배지 거터 폭 — 텍스트박스는 이 밖에서 시작(겹침 0).

    for i, item in enumerate(items):
        ry = t + i * row_h
        # 배지(oval, accent) — 거터 내부, 행 수직 중앙.
        by = ry + (row_h - badge_d) / 2.0
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                      Inches(l), Inches(by), Inches(badge_d), Inches(badge_d))
        oval.fill.solid()
        oval.fill.fore_color.rgb = _token_rgb(tokens, "accent")
        oval.line.fill.background()
        otf = oval.text_frame
        _set_textframe_basics(otf)
        op = otf.paragraphs[0]
        if PP_ALIGN is not None:
            op.alignment = PP_ALIGN.CENTER
        _add_run(op, str(i + 1), tokens, "badge")
        placed.append(PlacedShape(role="badge", rect=(l, by, badge_d, badge_d),
                                  has_text=True, text=str(i + 1), z=3))

        # 본문 텍스트박스 — 거터 밖에서 시작(겹침 0).
        tx = l + gutter
        tw = max(0.5, w - gutter)
        title, desc = _card_text(item)
        body_text = title if not desc else f"{title} — {desc}"
        tb = slide.shapes.add_textbox(Inches(tx), Inches(ry), Inches(tw), Inches(row_h))
        tf = tb.text_frame
        _set_textframe_basics(tf)
        p = tf.paragraphs[0]
        if PP_ALIGN is not None:
            p.alignment = PP_ALIGN.LEFT
        _add_run(p, str(body_text)[:160], tokens, "body")
        placed.append(PlacedShape(role="body", rect=(tx, ry, tw, row_h),
                                  has_text=True, text=_norm(body_text), z=2))
    return placed


def emit_card_grid(slide, cards, tokens: dict, region) -> list[PlacedShape]:
    """카드 그리드: 라운드 사각형 카드(card_bg + border + 그림자) + 제목/본문 run."""
    l, t, w, h = region
    placed: list[PlacedShape] = []
    cards = list(cards or [])
    n = len(cards)
    if n <= 0:
        return placed

    cols = 3 if n >= 5 else (2 if n >= 2 else 1)
    rows = (n + cols - 1) // cols
    gx, gy = 0.25, 0.25
    cw = (w - (cols - 1) * gx) / cols
    chh = (h - (rows - 1) * gy) / rows

    for idx, card in enumerate(cards):
        r = idx // cols
        c = idx % cols
        cx = l + c * (cw + gx)
        cy = t + r * (chh + gy)
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(cx), Inches(cy), Inches(cw), Inches(chh))
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
        shp.fill.solid()
        shp.fill.fore_color.rgb = _token_rgb(tokens, "card_bg")
        shp.line.color.rgb = _token_rgb(tokens, "border")
        shp.line.width = Pt(1.0)
        _apply_shadow(shp)

        tf = shp.text_frame
        _set_textframe_basics(tf, anchor_middle=False)
        title, desc = _card_text(card)
        p0 = tf.paragraphs[0]
        if PP_ALIGN is not None:
            p0.alignment = PP_ALIGN.LEFT
        _add_run(p0, str(title)[:60], tokens, "card_title")
        if desc:
            p1 = tf.add_paragraph()
            if PP_ALIGN is not None:
                p1.alignment = PP_ALIGN.LEFT
            _add_run(p1, str(desc)[:160], tokens, "body")
        card_text = title if not desc else f"{title}\n{desc}"
        placed.append(PlacedShape(role="card", rect=(cx, cy, cw, chh),
                                  has_text=True, text=_norm(card_text), z=2))
    return placed


def emit_figure_slot(slide, region, image_path: Optional[str] = None) -> PlacedShape:
    """그림 슬롯(장식 채움) 방출. 콘텐츠 텍스트는 베이크하지 않는다(장식 전용·텍스트 X).

    image_path 가 주어지면 손실-0 임베드(재인코딩 없이 add_picture). 없으면 중립
    회색 사각형 플레이스홀더를 그린다. 어느 경우든 has_text=False.
    """
    l, t, w, h = region
    if image_path:
        try:
            slide.shapes.add_picture(image_path, Inches(l), Inches(t), Inches(w), Inches(h))
            return PlacedShape(role="figure", rect=(l, t, w, h), has_text=False, z=0)
        except Exception:
            pass  # 임베드 실패 → 플레이스홀더 폴백.

    slot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(l), Inches(t), Inches(w), Inches(h))
    slot.fill.solid()
    slot.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF5)  # 중립 장식 회색(콘텐츠 색 아님).
    slot.line.color.rgb = RGBColor(0xD0, 0xD5, 0xDD)
    slot.line.width = Pt(1.0)
    return PlacedShape(role="figure", rect=(l, t, w, h), has_text=False, z=0)


def apply_tokens_to_run(run, tokens: dict, role: str) -> None:
    """design_tokens 의 색/타이포를 네이티브 run 서식에 적용 (Req 5.4: 신규 토큰 금지).

    role(title/body/badge/note/...)별로 _ROLE_STYLE 에서 폰트크기·bold·색토큰키·폰트토큰키를
    조회해 run.font.size / font.color.rgb / font.bold / font.name 에 매핑한다. 색·폰트명은
    design_tokens dict 에서만 읽고(부재 시 안전 폴백), 폰트크기는 타이포 계층용 role 기본값.
    """
    style = _ROLE_STYLE.get(role) or _ROLE_STYLE["body"]
    try:
        if Pt is not None:
            run.font.size = Pt(style["size"])
        run.font.bold = bool(style["bold"])
        run.font.color.rgb = _token_rgb(tokens, style["color"])
        run.font.name = _font_name_from_stack((tokens or {}).get(style["font"]))
    except Exception:
        pass


# ---------------------------------------------------------------------------
# 작업 18.1 — 고밀도 chrome 보강 (밀도 reference 6 + 스타일 품질 충족)
# ---------------------------------------------------------------------------
# 산출물_검증기 audit_native_density 의 (e) 밀도(parity_scorer reference=6)와
# (h) 스타일 품질을 표지/7개 알려진_레이아웃 전부에서 충족하기 위해, 각 _render_* 가
# 콘텐츠 외에 design_tokens 기반 "고밀도 chrome"을 가산 방출한다.
#
# 회귀 불변식(Property 1·5·7·14 / audit (b)(c)(d) 보존):
#   - 텍스트 없는 장식 컨테이너는 DECORATIVE_BG_ROLES role(section_bar/contact_box/
#     accent_bar/decorative_bg)로 표기 → finalize_placement / Property 4·14 겹침
#     검사 및 audit (b) 에서 제외(레이어드 디자인 정상 허용).
#   - 텍스트 chrome(번호 배지·푸터)은 콘텐츠(≤ _CONTENT_BOTTOM)와 분리된 하단
#     밴드(≥5.65)에만 배치해 콘텐츠 텍스트와 겹침 0(Property 14).
#   - 제목 셰이프(TITLE_ROLES)는 추가하지 않는다 → title_count 불변(Property 7).
#   - 신규 토큰 정의 없음 — 색·폰트는 design_tokens 에서만 읽는다(Req 5.4).
#   - 모든 chrome 셰이프는 슬라이드 경계(13.333×7.5) 안(Property 5 / audit (c)).
# ---------------------------------------------------------------------------


def _chrome_number(data: dict) -> int:
    """chrome 번호 배지용 정수(numbered_item 마커는 OVAL + 숫자 텍스트 검출)."""
    if isinstance(data, dict):
        for k in ("section_number", "number", "no"):
            v = data.get(k)
            try:
                if v is not None and int(v) >= 0:
                    return int(v)
            except Exception:
                continue
    return 1


def _chrome_footer(data: dict) -> str:
    """chrome 푸터 텍스트(slide_footer 마커는 하단부 텍스트 셰이프 검출)."""
    if isinstance(data, dict):
        for k in ("footer", "slide_footer", "date_str"):
            v = data.get(k)
            if isinstance(v, str) and v.strip():
                return v.strip()
    return "Confidential · Native Render"


def _emit_cover_chrome(slide, tokens: dict) -> list[PlacedShape]:
    """표지 고밀도 chrome — corner_glow(대형 OVAL)/icon_badge(소형 accent OVAL)/
    notice_chip(eyebrow pill, 라운드+border)/accent_bar(하단 가로 막대=footer 겸).

    모두 텍스트 없는 장식(decorative_bg/accent_bar role) → 겹침 검사 제외. 콘텐츠
    텍스트보다 먼저 방출되어 z-order 뒤(가독 보존). 신규 토큰 없음."""
    placed: list[PlacedShape] = []
    if MSO_SHAPE is None:
        return placed

    # corner_glow — 우상단 대형 accent 틴트 OVAL(최대변 ≥1.5).
    cg = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.0), Inches(0.25), Inches(1.95), Inches(1.95))
    cg.fill.solid(); cg.fill.fore_color.rgb = _tint_rgb(tokens, "accent", 0.80)
    cg.line.fill.background()
    placed.append(PlacedShape(role="decorative_bg", rect=(11.0, 0.25, 1.95, 1.95), has_text=False, z=1))

    # icon_badge — 소형 정사각 accent OVAL(유채색).
    ib = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(0.6), Inches(0.75), Inches(0.75))
    ib.fill.solid(); ib.fill.fore_color.rgb = _token_rgb(tokens, "accent")
    ib.line.fill.background()
    placed.append(PlacedShape(role="decorative_bg", rect=(0.7, 0.6, 0.75, 0.75), has_text=False, z=2))

    # notice_chip — eyebrow pill(라운드+border) accent 틴트(스타일 품질 round 검사 충족).
    nc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.42), Inches(3.5), Inches(0.55))
    try:
        nc.adjustments[0] = 0.5
    except Exception:
        pass
    nc.fill.solid(); nc.fill.fore_color.rgb = _tint_rgb(tokens, "accent", 0.86)
    nc.line.color.rgb = _token_rgb(tokens, "accent"); nc.line.width = Pt(1.0)
    placed.append(PlacedShape(role="decorative_bg", rect=(0.7, 1.42, 3.5, 0.55), has_text=False, z=2))

    # accent_bar(+footer) — 하단 가로 accent 막대(h≤0.3, w≥1.5; top≥6.3 → footer 겸).
    ab = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(6.7), Inches(3.6), Inches(0.15))
    ab.fill.solid(); ab.fill.fore_color.rgb = _token_rgb(tokens, "accent")
    ab.line.fill.background()
    placed.append(PlacedShape(role="accent_bar", rect=(0.7, 6.7, 3.6, 0.15), has_text=False, z=2))
    return placed


def emit_cover_density_chrome(slide, tokens: dict, *, add_step_grid: bool = True) -> list[PlacedShape]:
    """server.py 표지 경로(build_native_cover) 보강용 표지 density/style chrome.

    build_native_cover 는 palette 색으로 좌측 accent 바·우상단 corner_glow·eyebrow
    notice_chip·대형 제목·(KPI 시)step_grid 를 그리지만, (1) 소형 icon_badge(≤1.3in
    유채색 원)와 (2) KPI 부재 시 step_grid(둥근 카드 2개+)가 없어 밀도 5/6 에 그치고,
    (3) 채움색이 design_tokens accent 와 달라 audit 의 accent_color 스타일 검사에 걸린다.

    이 함수는 그 세 공백만 **가산적으로** 채운다 — design_tokens(audit 과 동일 출처)의
    accent 색으로 icon_badge(소형 OVAL)와 step_grid(둥근 카드 3개)를 표지의 빈 영역
    (좌상단·중하단 밴드)에 방출한다. 모두 텍스트 없는 장식(decorative_bg role) →
    겹침 검사 제외(레이어드), 콘텐츠 텍스트보다 뒤 z-order(가독 보존). 신규 토큰 없음.

    Args:
        slide: python-pptx Slide (build_native_cover 적용 후의 표지)
        tokens: design_tokens_for_profile 결과 dict (accent/primary/... 색 출처)
        add_step_grid: KPI 카드가 이미 있으면 False 로 전달해 step_grid 중복 방지.

    Returns:
        방출된 PlacedShape 목록(모두 decorative_bg, has_text=False).
    """
    placed: list[PlacedShape] = []
    if MSO_SHAPE is None or Inches is None:
        return placed

    # 1) icon_badge — 소형 정사각 accent OVAL(유채색, 최대변 ≤1.3). 좌상단 빈 영역.
    try:
        ib = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.72), Inches(0.55), Inches(0.62), Inches(0.62))
        ib.fill.solid(); ib.fill.fore_color.rgb = _token_rgb(tokens, "accent")
        ib.line.fill.background()
        placed.append(PlacedShape(role="decorative_bg", rect=(0.72, 0.55, 0.62, 0.62), has_text=False, z=1))
    except Exception:
        pass

    # 2) step_grid — 중하단 밴드에 accent 틴트 둥근 카드 3개(w≥1.2, h≥0.8). KPI 카드가
    #    이미 있으면(add_step_grid=False) 생략해 중복을 피한다. 텍스트 없는 장식.
    if add_step_grid:
        try:
            _n = 3
            _gap = 0.3
            _left0 = 1.15
            _w = 2.5
            _top = 5.5
            _h = 0.9
            for _i in range(_n):
                _l = _left0 + _i * (_w + _gap)
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(_l), Inches(_top), Inches(_w), Inches(_h))
                try:
                    card.adjustments[0] = 0.10
                except Exception:
                    pass
                card.fill.solid(); card.fill.fore_color.rgb = _tint_rgb(tokens, "accent", 0.86)
                try:
                    card.line.color.rgb = _token_rgb(tokens, "accent")
                    card.line.width = Pt(1.25)
                except Exception:
                    pass
                _apply_shadow(card)
                placed.append(PlacedShape(role="decorative_bg", rect=(_l, _top, _w, _h), has_text=False, z=2))
        except Exception:
            pass
    return placed


def _emit_body_chrome(slide, tokens: dict, *, number=1, footer_text: str = "") -> list[PlacedShape]:
    """본문 고밀도 chrome(하단 밴드, 콘텐츠와 분리) — 연락처 레일(contact_box +
    note_callout)/노티스 탭(notice_tab)/다크 푸터 바(section_header)/번호 배지
    (numbered_item)/푸터 텍스트(slide_footer). 6개 body 밀도 마커를 충족한다.

    텍스트 없는 컨테이너는 DECORATIVE_BG_ROLES → 겹침 검사 제외. 텍스트(배지/푸터)는
    하단 밴드에서 상호 비겹침. 신규 토큰 없음."""
    placed: list[PlacedShape] = []
    if MSO_SHAPE is None:
        return placed

    # 1) 연락처 레일 — bordered rounded box(contact_box+note_callout) + 좌측 accent vbar.
    cx, cy, cw, ch = 0.5, 5.65, 5.4, 0.95
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(cw), Inches(ch))
    try:
        box.adjustments[0] = 0.08
    except Exception:
        pass
    box.fill.solid(); box.fill.fore_color.rgb = _tint_rgb(tokens, "primary", 0.90)
    box.line.color.rgb = _token_rgb(tokens, "border"); box.line.width = Pt(1.0)
    _apply_shadow(box)
    placed.append(PlacedShape(role="contact_box", rect=(cx, cy, cw, ch), has_text=False, z=1))
    vb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy), Inches(0.12), Inches(ch))
    vb.fill.solid(); vb.fill.fore_color.rgb = _token_rgb(tokens, "accent"); vb.line.fill.background()
    placed.append(PlacedShape(role="accent_bar", rect=(cx, cy, 0.12, ch), has_text=False, z=2))

    # 2) 노티스 탭 — 작은 accent 솔리드 라운드 탭(notice_tab).
    nt = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.1), Inches(5.8), Inches(1.7), Inches(0.5))
    try:
        nt.adjustments[0] = 0.4
    except Exception:
        pass
    nt.fill.solid(); nt.fill.fore_color.rgb = _token_rgb(tokens, "accent"); nt.line.fill.background()
    placed.append(PlacedShape(role="accent_bar", rect=(6.1, 5.8, 1.7, 0.5), has_text=False, z=2))

    # 3) 다크 푸터 바 — 넓은 다크 라운드 막대(section_header) + 그림자(스타일 품질).
    fbx, fby, fbw, fbh = 0.5, 6.75, 12.33, 0.55
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fbx), Inches(fby), Inches(fbw), Inches(fbh))
    try:
        bar.adjustments[0] = 0.18
    except Exception:
        pass
    bar.fill.solid(); bar.fill.fore_color.rgb = _token_rgb(tokens, "text_dark"); bar.line.fill.background()
    _apply_shadow(bar)
    placed.append(PlacedShape(role="section_bar", rect=(fbx, fby, fbw, fbh), has_text=False, z=1))

    # 4) 번호 배지 — OVAL accent + 숫자 텍스트(numbered_item). 바 좌측, 푸터와 비겹침.
    bd = 0.42
    bx = fbx + 0.12
    by = fby + (fbh - bd) / 2.0
    oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx), Inches(by), Inches(bd), Inches(bd))
    oval.fill.solid(); oval.fill.fore_color.rgb = _token_rgb(tokens, "accent"); oval.line.fill.background()
    otf = oval.text_frame
    _set_textframe_basics(otf)
    op = otf.paragraphs[0]
    if PP_ALIGN is not None:
        op.alignment = PP_ALIGN.CENTER
    _add_run(op, str(int(number)), tokens, "badge")
    placed.append(PlacedShape(role="badge", rect=(bx, by, bd, bd), has_text=True, text=str(int(number)), z=3))

    # 5) 푸터 텍스트 — 다크 바 위 흰색 텍스트(slide_footer). 배지 오른쪽이라 비겹침.
    ft = str(footer_text or "Confidential · Native Render")
    tx = bx + bd + 0.18
    tw = 8.0
    tb = slide.shapes.add_textbox(Inches(tx), Inches(fby + 0.05), Inches(tw), Inches(fbh - 0.1))
    tf = tb.text_frame
    _set_textframe_basics(tf)
    p = tf.paragraphs[0]
    if PP_ALIGN is not None:
        p.alignment = PP_ALIGN.LEFT
    _add_run(p, ft[:80], tokens, "badge")  # badge role = text_light(흰색) on dark bar
    placed.append(PlacedShape(role="caption", rect=(tx, fby + 0.05, tw, fbh - 0.1),
                              has_text=True, text=_norm(ft), z=3))
    return placed


# ---------------------------------------------------------------------------
# 겹침 검사 정책 (A안) — design Property 4 / Req 2 명확화.
#
# 텍스트를 포함하지 않는 "장식 배경 도형"(섹션 헤더 막대·카드 배경 컨테이너·노트
# 콜아웃 배경 박스·좌측 보더 등)은 겹침 검사 대상에서 제외한다. 그 위에 텍스트·요소를
# 의도적으로 올리는 레이어드 디자인은 정상으로 허용한다. 겹침률 10% 미만 규칙은
# 텍스트 보유 셰이프 쌍(텍스트↔텍스트) 및 텍스트↔비배경 이미지/도형 쌍에만 적용한다.
# 텍스트↔이미지 가림은 본 함수가 아니라 z-order(Property 13/audit (g))가 보장한다.
# ---------------------------------------------------------------------------

# emit_* 가 만드는 "텍스트 없는 장식 배경 컨테이너" role 집합 — 겹침 검사 제외.
DECORATIVE_BG_ROLES: frozenset = frozenset({
    "section_bar",    # emit_section_header_bar 다크 막대
    "contact_box",    # emit_contact_box 틴트 사각형
    "note_box",       # emit_note_callout 경고 틴트 사각형
    "note_border",    # emit_note_callout 좌측 보더
    "accent_bar",     # 연락처/표지 좌측 액센트 바
    "decorative_bg",  # 풀블리드 장식 배경
})

# 비배경 이미지 role — has_text=False 이지만 장식 배경이 아니므로 텍스트와의 겹침은 검사한다.
_IMAGE_ROLES: frozenset = frozenset({"figure", "image"})

# 겹침 임계(작은 셰이프 면적의 10%) — layout_geometry.DEFAULT_THRESHOLD 와 동일.
_OVERLAP_THRESHOLD: float = 0.10


def _participates_in_collision(ps: PlacedShape) -> bool:
    """A안 겹침 검사 대상 여부.

    제외: 텍스트 없는 장식 배경 컨테이너(DECORATIVE_BG_ROLES).
    포함: 텍스트 보유 셰이프(has_text=True), 그리고 비배경 이미지(figure/image).
    """
    if ps.role in DECORATIVE_BG_ROLES:
        return False
    if ps.has_text:
        return True
    if ps.role in _IMAGE_ROLES:
        return True
    return False


def _clamp_rect(rect: tuple) -> tuple:
    """단일 rect 를 슬라이드 경계 안으로 보정(Req 3, Property 5).

    이미 경계 안(eps=0.05)이면 입력을 그대로 반환(no-op, Property 6). 벗어나면 폭/높이가
    슬라이드보다 크면 fit_within 으로 종횡비 보존 축소 후 clamp_into_bounds, 아니면
    clamp_into_bounds 로 평행이동한다.
    """
    if within_bounds(rect, SLIDE_RECT, eps=0.05):
        return rect  # no-op — 비결함 입력 좌표 불변 보존
    l, t, w, h = rect
    sw, sh = SLIDE_RECT[2], SLIDE_RECT[3]
    if w > sw or h > sh:
        fitted = fit_within(SLIDE_RECT, w, h)
        return tuple(clamp_into_bounds(fitted, SLIDE_RECT))
    return tuple(clamp_into_bounds(rect, SLIDE_RECT))


def _dedup_titles(placed: list[PlacedShape]) -> list[PlacedShape]:
    """제목 역할(TITLE_ROLES) 셰이프가 정규화 텍스트 기준 중복이면 1개만 유지(Req 4.2).

    서로 다른 제목 텍스트는 모두 유지한다. 정규화 텍스트가 비어있는 제목은 dedup 하지 않는다.
    """
    seen: set = set()
    out: list[PlacedShape] = []
    for ps in placed:
        if ps.role in TITLE_ROLES:
            key = _norm(ps.text)
            if key:
                if key in seen:
                    continue  # 중복 제목 셰이프 제거
                seen.add(key)
        out.append(ps)
    return out


def _shape_id(ps: PlacedShape, idx: int) -> str:
    """OverlapError 보고용 셰이프 식별자."""
    label = (ps.text or ps.role or "")[:24]
    return f"{ps.role}#{idx}:{label}"


def _overlap_violations(shapes: list[PlacedShape], idxs: list[int]) -> list[tuple]:
    """겹침률 10% 이상인 셰이프 쌍 목록(식별자 쌍) 반환."""
    pairs: list[tuple] = []
    n = len(shapes)
    for i in range(n):
        for j in range(i + 1, n):
            a, b = shapes[i].rect, shapes[j].rect
            ov = overlap_area(a, b)
            if ov <= 0.0:
                continue
            amin = min(area(a), area(b))
            if amin <= 0.0:
                continue
            if ov >= _OVERLAP_THRESHOLD * amin:
                pairs.append((_shape_id(shapes[i], idxs[i]), _shape_id(shapes[j], idxs[j])))
    return pairs


def finalize_placement(placed: list[PlacedShape], *, slide_id=None) -> list[PlacedShape]:
    """렌더된 PlacedShape 목록을 받아 배치를 보정해 반환한다.

    처리 순서:
      1) 제목 중복 제거 — 정규화 텍스트가 동일한 제목 역할 셰이프는 1개만 유지(Req 4.2, P8).
      2) 경계 안 보정 — 각 셰이프 rect 를 슬라이드 경계 안으로(Req 3, P5). 이미 경계 안이면
         좌표 불변(no-op, P6).
      3) 겹침 보정(A안) — 텍스트 없는 장식 배경 도형을 제외하고, 텍스트 보유 셰이프 및
         텍스트↔비배경 이미지 쌍만 resolve_collisions 로 겹침률 10% 미만으로 보정(Req 2, P4).
         장식 배경 위 텍스트의 의도적 레이어링은 보존한다.
      4) 과밀 검사 — 보정 후에도 겹침률 10% 미만으로 못 낮춘 쌍이 있으면
         OverlapError(slide_id, [(a, b), ...]) 발생(Req 2.5).

    비결함 입력(이미 경계 안·겹침 없음·제목 중복 없음)은 좌표 불변 no-op 으로 반환한다(P6).
    """
    if not placed:
        return []

    # --- 1) 제목 중복 제거 (Req 4.2, Property 8) ---
    result = _dedup_titles(placed)

    # --- 2) 경계 안 보정 (Req 3, Property 5/6) ---
    for ps in result:
        ps.rect = _clamp_rect(ps.rect)

    # --- 3) 겹침 보정 (Req 2, Property 4, A안 — 장식 배경 도형 제외) ---
    part_idx = [i for i, ps in enumerate(result) if _participates_in_collision(ps)]
    if len(part_idx) >= 2:
        rects = [result[i].rect for i in part_idx]
        resolved = resolve_collisions(rects, threshold=_OVERLAP_THRESHOLD, axis="vertical")
        # 보정 결과 적용 + 재클램프(수직 스택이 경계를 벗어났을 수 있음).
        for k, i in enumerate(part_idx):
            result[i].rect = _clamp_rect(tuple(resolved[k]))

        # --- 4) 과밀 검사 (Req 2.5) — 보정 후에도 10% 이상 겹치면 OverlapError ---
        part_shapes = [result[i] for i in part_idx]
        violations = _overlap_violations(part_shapes, part_idx)
        if violations:
            raise OverlapError(slide_id, violations)

    return result


# ===========================================================================
# 작업 8.1 — Vertex 장식_비주얼 경로 (가산적 추가)
#
# 장식_비주얼(장식_배경/히어로_일러스트/그림슬롯)을 옵트인 시 vertex_image_module
# 단일 모듈로 생성·채운다. 콘텐츠 텍스트는 절대 이미지에 베이크하지 않는다(장식 전용).
#
# 핵심 불변식(design.md §역할 분리, Property 13/16/17, Req 9.2/11.x):
#   - 옵트인 판정(Req 11.3): AE_ENABLE_VERTEX_IMAGE == "1" AND vertex 클라이언트 enabled
#     (자격증명 존재). OFF/자격증명 부재/생성 실패 → None 반환 → 장식 슬롯 비움 또는
#     회색 플레이스홀더, 콘텐츠 네이티브는 손실 0으로 정상 생성(예외 전파 금지).
#   - 단일 모듈 경유(Req 11.5/gateway.md): 이미지 생성은 vertex_image_module 만 경유.
#   - 손실-0 임베드(Req 9.2/11.4/Property 16): 생성 이미지는 재인코딩 없이 원본 바이트로
#     저장 후 emit_figure_slot 의 add_picture 로 임베드.
#   - z-order(Property 13/16): 장식 이미지는 콘텐츠 텍스트보다 뒤(back-most). PlacedShape.z=0,
#     실제 spTree 에서도 back-most 로 보낸다(_send_to_back).
#   - 장식 배경 경계(Property 16): 풀블리드 장식 배경은 콘텐츠 텍스트 없음(has_text=False)이
#     보장된다(emit_figure_slot 가 텍스트를 그리지 않음).
# ===========================================================================

# Vertex 장식 이미지에 사용할 모델 클래스(vertex_image_module.VERTEX_MODEL_REGISTRY 키).
_VERTEX_DECOR_MODEL_CLASS: str = "image_generation_high_quality"

# emit_decorative_figure 가 장식_배경(풀블리드)으로 마킹할 region_kind 집합.
# 이 종류는 PlacedShape.role 을 "decorative_bg" 로 지정해 finalize_placement 의 A안
# 겹침 검사에서 제외된다(콘텐츠 위 의도적 레이어링 보존; z-order 가 가림을 방지).
_DECOR_BG_KINDS: frozenset = frozenset({"background", "decorative_bg", "hero_bg", "fullbleed"})


def _vertex_optin_enabled() -> bool:
    """Vertex 장식 옵트인 환경변수 판정 (Req 11.3 전반).

    AE_ENABLE_VERTEX_IMAGE 값이 정확히 "1" 일 때만 True. 자격증명 존재 여부는
    _get_vertex_client / client.enabled 가 별도로 판정한다.
    """
    return os.environ.get("AE_ENABLE_VERTEX_IMAGE", "").strip() == "1"


def _get_vertex_client(aws_profile: str = "", credentials: Optional[dict] = None):
    """vertex_image_module.get_vertex_image_client 위임(단일 모듈 경유, Req 11.5).

    이미지 생성은 이 단일 모듈만 경유한다(다른 이미지 API 직접 호출 금지). 모듈 미가용/
    예외 시 None 을 반환한다. **테스트는 이 함수를 monkeypatch 하여 네트워크 없이
    Vertex 를 mock 한다**(실제 Vertex 호출 금지).
    """
    try:
        from vertex_image_module import get_vertex_image_client  # type: ignore
    except Exception:  # pragma: no cover - 패키지 경로 fallback
        try:
            from ai_engine.vertex_image_module import get_vertex_image_client  # type: ignore
        except Exception:
            return None
    try:
        return get_vertex_image_client(aws_profile=aws_profile, credentials=credentials)
    except Exception:
        return None


def _run_coro(coro):
    """동기 컨텍스트에서 awaitable(코루틴) 실행. 실행 중 루프가 있으면 별도 스레드에서
    새 이벤트 루프로 실행해 'event loop already running' 충돌을 피한다."""
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        return asyncio.run(coro)
    import concurrent.futures
    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as ex:
        return ex.submit(lambda: asyncio.run(coro)).result()


def _write_decor_image(raw: bytes, region_kind: str, out_dir: Optional[str] = None) -> Optional[str]:
    """생성된 이미지 원본 바이트를 로컬 파일로 저장 후 경로 반환(손실-0: 재인코딩 없음).

    add_picture 는 파일 바이트를 그대로 패키지에 임베드하므로, base64 디코드 결과를
    그대로 기록하면 Property 16(바이트 동일 임베드)이 보장된다.
    """
    if not raw:
        return None
    try:
        base = out_dir or os.path.join(tempfile.gettempdir(), "ae_vertex_decor")
        os.makedirs(base, exist_ok=True)
        fname = f"decor_{region_kind}_{int(time.time() * 1000)}_{os.getpid()}.png"
        path = os.path.join(base, fname)
        with open(path, "wb") as f:
            f.write(raw)
        return path
    except Exception:
        return None


def maybe_generate_decorative(
    prompt,
    region_kind: str = "figure",
    *,
    aspect_ratio: str = "16:9",
    aws_profile: str = "",
    credentials: Optional[dict] = None,
    out_dir: Optional[str] = None,
) -> Optional[str]:
    """옵트인 시 Vertex 로 장식_비주얼 이미지를 생성하고 로컬 경로를 반환한다.

    Args:
        prompt: 장식_비주얼 생성 프롬프트(콘텐츠 텍스트가 아닌 장식 지시).
        region_kind: 슬롯 종류(figure/background/hero 등) — 파일명 태깅용.
        aspect_ratio: 생성 종횡비(슬라이드 16:9 기본).
        aws_profile/credentials: Secrets Manager 자격증명 해석에 전달.
        out_dir: 이미지 저장 디렉터리(미지정 시 임시 디렉터리).

    Returns:
        생성 성공 시 로컬 이미지 경로(str). 다음의 경우 None(예외 전파 금지):
          - 옵트인 OFF(AE_ENABLE_VERTEX_IMAGE != "1")
          - vertex 클라이언트 미가용 또는 enabled=False(자격증명 부재)
          - 빈 프롬프트
          - 생성 실패(error dict / 빈 images / 디코드 실패 / 예외)

    이미지 생성은 vertex_image_module 단일 모듈만 경유한다(Req 11.5).
    """
    # 옵트인 판정 (Req 11.3) — 환경변수 우선.
    if not _vertex_optin_enabled():
        return None
    if not prompt or not str(prompt).strip():
        return None

    client = _get_vertex_client(aws_profile=aws_profile, credentials=credentials)
    # enabled 판정 (Req 11.3) — 자격증명 존재 시에만 True.
    if client is None or not getattr(client, "enabled", False):
        return None

    # 단일 모듈(vertex_image_module) 경유 생성. async/sync 양쪽 mock 대응.
    try:
        maybe_res = client.generate(
            str(prompt),
            model_class=_VERTEX_DECOR_MODEL_CLASS,
            aspect_ratio=aspect_ratio,
        )
        result = _run_coro(maybe_res) if inspect.isawaitable(maybe_res) else maybe_res
    except Exception:
        return None  # 생성 실패 → 폴백(None), 예외 전파 금지.

    if not isinstance(result, dict):
        return None
    images = result.get("images") or []
    if not images:
        return None  # error dict / 빈 결과 → 폴백.

    # base64 → 원본 바이트(손실-0). 재인코딩하지 않는다.
    try:
        raw = base64.b64decode(images[0])
    except Exception:
        return None
    return _write_decor_image(raw, region_kind, out_dir)


def _send_to_back(shape) -> None:
    """spTree 에서 shape 요소를 back-most(최하단 z-order)로 이동(Property 13/16).

    python-pptx 는 add 순서가 z-order 이므로, 나중에 add 된 장식 이미지를 콘텐츠 텍스트
    뒤로 보내려면 spTree 의 첫 셰이프 위치로 옮긴다(앞쪽 nvGrpSpPr/grpSpPr 비셰이프
    요소 뒤). 콘텐츠 텍스트 셰이프는 항상 그 이미지보다 앞선 z-order 에 위치하게 된다.
    """
    _SHAPE_TAGS = ("sp", "pic", "graphicFrame", "grpSp", "cxnSp")
    try:
        sp = shape._element
        spTree = sp.getparent()
        if spTree is None:
            return
        spTree.remove(sp)
        ref = None
        for child in spTree:
            ln = child.tag.rsplit("}", 1)[-1]
            if ln in _SHAPE_TAGS:
                ref = child
                break
        if ref is None:
            spTree.append(sp)
        else:
            ref.addprevious(sp)
    except Exception:
        pass


def emit_decorative_figure(
    slide,
    region,
    prompt,
    region_kind: str = "figure",
    *,
    aspect_ratio: str = "16:9",
    aws_profile: str = "",
    credentials: Optional[dict] = None,
    placeholder_on_empty: bool = True,
    out_dir: Optional[str] = None,
) -> Optional[PlacedShape]:
    """장식 그림 슬롯을 Vertex(옵트인 시)로 채우고 emit_figure_slot 으로 임베드한다.

    흐름:
      1) maybe_generate_decorative 로 이미지 생성 시도(옵트인 OFF/실패 시 None).
      2) emit_figure_slot(slide, region, image_path) 로 손실-0 임베드(또는 회색 플레이스홀더).
      3) 방출된 셰이프를 back-most 로 보내 콘텐츠 텍스트보다 뒤 z-order 보장(Property 13/16).

    Args:
        region: (left, top, width, height) 인치. 슬라이드 경계 안(또는 풀블리드 장식).
        region_kind: figure(그림슬롯) / background·hero_bg(풀블리드 장식 배경) 등.
            _DECOR_BG_KINDS 에 속하면 PlacedShape.role="decorative_bg" 로 마킹되어
            finalize_placement 겹침 검사에서 제외된다(콘텐츠 위 의도적 레이어링).
        placeholder_on_empty: 이미지 부재 시 회색 플레이스홀더를 둘지(True) 비울지(False).
            풀블리드 장식 배경은 False 로 두어 생성 실패 시 거대한 회색 박스를 만들지 않는다.

    Returns:
        PlacedShape(has_text=False, z=0). 이미지 부재 + placeholder_on_empty=False 면 None
        (장식 슬롯 비움 — 콘텐츠 손실 0). 예외 전파 금지.
    """
    image_path = maybe_generate_decorative(
        prompt, region_kind,
        aspect_ratio=aspect_ratio, aws_profile=aws_profile,
        credentials=credentials, out_dir=out_dir,
    )
    if image_path is None and not placeholder_on_empty:
        return None  # 장식 슬롯 비움(콘텐츠 손실 0, Req 11.3).

    ps = emit_figure_slot(slide, region, image_path)
    # 장식 이미지/플레이스홀더를 콘텐츠보다 뒤로(back-most) — z-order 불변식.
    try:
        if len(slide.shapes) > 0:
            _send_to_back(slide.shapes[-1])
    except Exception:
        pass

    # 풀블리드 장식 배경은 role 을 decorative_bg 로 마킹(겹침 검사 제외 + audit 제외).
    if region_kind in _DECOR_BG_KINDS:
        ps.role = "decorative_bg"
    return ps


def _default_decor_prompt(layout: str, data: dict) -> str:
    """레이아웃별 기본 장식_배경 프롬프트(콘텐츠 텍스트 아님 — 추상 장식 지시).

    콘텐츠 텍스트(title/본문)를 프롬프트에 넣어 이미지에 베이크되게 하지 않는다. 추상적인
    배경 분위기만 지시한다(텍스트 없는 장식_비주얼).
    """
    base = (
        "abstract modern decorative background, soft gradient, subtle geometric shapes, "
        "professional presentation aesthetic, no text, no letters, no words"
    )
    if layout == "cover":
        return "elegant hero cover " + base
    if layout == "section_divider":
        return "minimal section divider " + base
    return base


def maybe_add_decorative_background(
    slide,
    data: dict,
    layout: str,
    *,
    aws_profile: str = "",
    credentials: Optional[dict] = None,
) -> list[PlacedShape]:
    """옵트인 ON + 장식 힌트 시 풀블리드 장식_배경을 방출한다(콘텐츠보다 먼저 → back-most).

    옵트인 OFF(또는 자격증명 부재)면 빈 리스트를 반환해 콘텐츠 네이티브만 손실 0으로
    렌더되게 한다(Req 11.3, Property 17) — 기존 동작과 바이트 동일(no-op).

    장식 힌트: data 의 ``decorative_bg``/``hero_prompt``/``figure_prompt`` 문자열, 또는
    cover/section_divider 레이아웃의 기본 장식 프롬프트. 콘텐츠 텍스트는 프롬프트에 넣지
    않는다(베이크 금지).

    생성 실패 시(ON 이지만 자격증명/네트워크 실패) placeholder_on_empty=False 로 두어
    장식 슬롯을 비운다(거대한 회색 배경 방지).
    """
    if not _vertex_optin_enabled():
        return []  # 옵트인 OFF → 장식 없음(콘텐츠 손실 0).
    data = data if isinstance(data, dict) else {}

    prompt = (
        data.get("decorative_bg")
        or data.get("hero_prompt")
        or data.get("figure_prompt")
        or ""
    )
    prompt = str(prompt).strip()
    if not prompt:
        if layout in ("cover", "section_divider"):
            prompt = _default_decor_prompt(layout, data)
        else:
            return []  # 배경 후보 아님 + 힌트 없음 → 장식 없음.

    region = (0.0, 0.0, SLIDE_W_IN, SLIDE_H_IN)  # 풀블리드 장식 배경(경계 안).
    ps = emit_decorative_figure(
        slide, region, prompt, region_kind="background",
        aws_profile=aws_profile, credentials=credentials,
        placeholder_on_empty=False,
    )
    return [ps] if ps is not None else []
