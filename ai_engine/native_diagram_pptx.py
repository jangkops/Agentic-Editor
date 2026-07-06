"""Native (editable) PPTX diagram builder.

문제 2 해결 — 슬라이드 안의 도표를 통짜 PNG(`add_picture`)가 아니라 PowerPoint에서
개별 요소 편집이 가능한 *네이티브 도형*(rounded rectangle + text frame + connector)으로
직접 조립한다. 젠스파크/감마 슬라이드처럼, 다운로드 후 도형·텍스트·화살표를 각각 수정할 수 있다.

설계:
- tree/flow/block 세 종류 파싱. 시각 결과가 PNG 버전과 의미적으로 일치한다.
- 0~100 가상 좌표계를 슬라이드의 지정 사각 영역(EMU)으로 선형 매핑한다.
- python-pptx만 사용 (Bedrock/Vertex 무관 — 순수 로컬 도형 조립).
- 어떤 입력에도 예외를 호출자로 전파하지 않는다(폴백 격리). 실패 시 False 반환.

좌표계:
- 가상 캔버스 (vx, vy) ∈ [0,100]×[0,100], 좌하단 원점.
- vy=100 → 영역 top, vy=0 → 영역 bottom (y축 뒤집기).
"""
from __future__ import annotations


def _hex_to_rgb(h: str):
    """'#RRGGBB' 또는 'RRGGBB' → (r,g,b). 실패 시 None."""
    if not isinstance(h, str):
        return None
    s = h.strip().lstrip("#")
    if len(s) != 6:
        return None
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except ValueError:
        return None


def _palette_colors(palette):
    """팔레트 리스트 → [hex, ...]. 기본값은 호출부에서 보강."""
    out = []
    if isinstance(palette, (list, tuple)):
        for c in palette:
            rgb = _hex_to_rgb(c)
            if rgb:
                out.append(c.strip().lstrip("#").upper())
    return out


def _parse_tree(content: str):
    """들여쓰기 깊이 파싱. → [{depth, text, is_dir}]."""
    lines_raw = [ln for ln in content.splitlines() if ln.strip()][:40]

    def _depth(s: str) -> int:
        i = 0
        d = 0
        while i < len(s):
            if s[i] == "\t":
                d += 1
                i += 1
            elif s.startswith("    ", i):
                d += 1
                i += 4
            elif s.startswith("  ", i):
                d += 1
                i += 2
            elif s[i] in "│ ":
                i += 1
            else:
                break
        return d

    nodes = []
    for raw in lines_raw:
        d = _depth(raw)
        text = raw
        for ch in ("│", "├", "└", "─", "—", "├──", "└──"):
            text = text.replace(ch, "")
        text = text.strip().lstrip("-*•").strip()
        if not text:
            continue
        is_dir = text.endswith("/") or text.endswith("\\")
        nodes.append({"depth": d, "text": text.rstrip("/\\"), "is_dir": is_dir})
    return nodes


def _parse_flow(content: str):
    """-> / → / 줄바꿈 분리. → [str]."""
    raw = content.replace("→", "->").replace("⇒", "->")
    parts = []
    for chunk in raw.split("\n"):
        for p in chunk.split("->"):
            t = p.strip().lstrip("-*•").strip()
            if t:
                parts.append(t[:48])
    return parts[:10]


def _parse_block(content: str):
    """줄바꿈 분리 항목. → [str]."""
    lines = [ln.strip().lstrip("-*•").strip() for ln in content.splitlines()]
    return [ln for ln in lines if ln][:10]


def _strip_md(s):
    """마크다운 인라인 마커 제거(**bold**, `code`, 선두 -/#/•). 네이티브 텍스트 공통."""
    import re as _re_sm
    t = str(s or "")
    t = _re_sm.sub(r"`([^`]+)`", r"\1", t)
    t = t.replace("**", "").replace("__", "")
    t = _re_sm.sub(r"^\s*(?:[-*\u2022]\s+|#{1,6}\s+)", "", t)
    return t


def _parse_cards(content: str):
    """카드 그리드 파싱 → [(title, desc), ...].

    각 줄을 카드 1개로. '제목: 설명' / '제목 - 설명' / '제목 | 설명' / '제목 — 설명'으로
    제목·설명 분리(구분자 없으면 설명 빈 문자열). 최대 6개.
    """
    out = []
    for ln in content.splitlines():
        s = ln.strip().lstrip("-*•").strip()
        if not s:
            continue
        title, desc = s, ""
        for sep in (" — ", " - ", " | ", ": ", "：", "—", "|"):
            if sep in s:
                left, right = s.split(sep, 1)
                title, desc = left.strip(), right.strip()
                break
        out.append((title[:44], desc[:100]))
        if len(out) >= 6:
            break
    return out


def _parse_kpis(content: str):
    """KPI/지표 라인 파싱 → [(value, label, delta), ...] (최대 4).

    지원 형식:
      "매출 증가: 32% (+8%p)" / "고객 만족도 | 95% | +3%" / "신규 고객 320명"
    """
    import re as _re
    out = []
    for ln in content.splitlines():
        s = ln.strip().lstrip("-*\u2022").strip()
        if not s:
            continue
        delta = ""
        m = _re.search(r"[(\uff08]\s*([+\-\u25b2\u25bc\u25b3\u25bd]?\s*[\d.]+\s*%?\s*p?)\s*[)\uff09]", s)
        if m:
            delta = m.group(1).replace(" ", "")
            s = (s[:m.start()] + s[m.end():]).strip()
        label, value = "", ""
        for sep in (" | ", "|", ": ", "\uff1a", " \u2014 ", " - "):
            if sep in s:
                parts = [p.strip() for p in s.split(sep) if p.strip()]
                if len(parts) >= 2:
                    label, value = parts[0], parts[1]
                    if not delta and len(parts) >= 3:
                        delta = parts[2]
                break
        if not value:
            vm = _re.search(r"([+\-]?[\d,]+(?:\.\d+)?\s*(?:%p|%|\uc5b5|\ub9cc|\uba85|\uac74|\uc810|\uac1c|\uc704|\ubc30|x|X)?)", s)
            if vm and any(c.isdigit() for c in vm.group(1)):
                value = vm.group(1).strip()
                label = (s[:vm.start()] + s[vm.end():]).strip(" :|-\u2014\u00b7")
            else:
                continue
        if not label:
            label, value = value, ""
        out.append((value[:14], label[:30], delta[:12]))
        if len(out) >= 4:
            break
    return out


def _parse_progress(content: str):
    """진행률/목표 달성 라인 파싱 → [(label, percent_int, status), ...] (최대 6)."""
    import re as _re
    _status_kw = ("완료", "진행중", "진행 중", "지연", "보류", "대기", "위험", "정상",
                  "done", "completed", "in progress", "ongoing", "at risk",
                  "delayed", "blocked", "on track")
    out = []
    for ln in content.splitlines():
        s = ln.strip().lstrip("-*\u2022").strip()
        if not s:
            continue
        pm = _re.search(r"(\d{1,3})\s*%", s)
        if not pm:
            continue
        pct = max(0, min(100, int(pm.group(1))))
        s2 = (s[:pm.start()] + s[pm.end():])
        status = ""
        low = s2.lower()
        for kw in _status_kw:
            idx = s2.find(kw)
            if idx < 0:
                idx = low.find(kw.lower())
            if idx >= 0:
                status = kw
                s2 = (s2[:idx] + s2[idx + len(kw):])
                break
        label = s2.strip(" :|-\u2014\u00b7()[]").strip()
        if not label:
            label = "항목"
        out.append((label[:34], pct, status[:10]))
        if len(out) >= 6:
            break
    return out


def _text_units(text: str) -> float:
    """텍스트의 시각적 폭 추정(단위: 'ASCII 대문자 1글자' 근사).

    한글/CJK/전각은 1.0, ASCII 영문·숫자는 0.6, 공백은 0.35, 기타는 0.6.
    박스 폭을 텍스트에 맞춰 산정해 잘림/과여백을 줄이는 데 쓴다.
    """
    u = 0.0
    for ch in str(text):
        o = ord(ch)
        if o >= 0x1100 and (0x1100 <= o <= 0x115F or 0x2E80 <= o <= 0xA4CF or
                            0xAC00 <= o <= 0xD7A3 or 0xF900 <= o <= 0xFAFF or
                            0xFE30 <= o <= 0xFE4F or 0xFF00 <= o <= 0xFF60 or
                            0xFFE0 <= o <= 0xFFE6):
            u += 1.0
        elif ch == " ":
            u += 0.35
        elif ch.isalnum():
            u += 0.6
        else:
            u += 0.55
    return u


# ── 의미 아이콘 매핑 — 카드 제목/설명 키워드 → 의미 키(Lucide 아이콘 + autoshape 폴백) ──
# 1순위: Lucide 고품질 벡터 아이콘(icon_assets, Chrome 렌더 PNG).
# 폴백: PowerPoint 네이티브 autoshape(_KEY_TO_MSO). 둘 다 칩 안 흰색.
_ICON_MAP = [
    (("속도", "성능", "빠른", "실시간", "real-time", "realtime", "fast", "speed",
      "performance", "즉시", "원클릭", "원스톱", "신속", "가속"), "zap"),
    (("클라우드", "cloud", "aws", "서버", "server", "인프라", "infra", "ec2",
      "인스턴스", "instance", "리전", "region"), "cloud"),
    (("설정", "운영", "프로세스", "process", "자동화", "automation", "ops", "gear",
      "config", "구성", "파이프라인", "pipeline", "프로비저닝", "provision",
      "온보딩", "onboarding", "관리"), "gear"),
    (("데이터", "저장소", "db", "database", "repository", "repo", "storage",
      "백업", "backup", "자산", "dataset", "데이터셋"), "database"),
    (("보안", "통제", "보호", "shield", "security", "protect", "감사", "audit",
      "무결성", "권한", "permission", "governance", "거버넌스", "정책", "policy",
      "차단", "방지"), "shield"),
    (("비용", "예산", "budget", "cost", "요금", "과금", "billing", "결제",
      "payment", "절감", "saving", "할인", "달러"), "dollar"),
    (("분석", "통계", "모니터링", "monitor", "analytics", "report", "대시보드",
      "dashboard", "지표", "metric", "kpi", "현황", "추이"), "chart"),
    (("문서", "가이드", "document", "doc", "guide", "파일", "리포트", "보고서",
      "매뉴얼", "manual", "계약", "contract", "제안"), "file"),
    (("증가", "성장", "상승", "growth", "increase", "향상", "개선", "확대",
      "expand", "scale"), "up"),
    (("감소", "하락", "reduce", "decrease", "축소", "downsize"), "down"),
    (("연결", "통합", "integration", "integrate", "connect", "연동", "링크",
      "link", "매핑", "mapping", "바인딩", "binding"), "refresh"),
    (("핵심", "중요", "key", "highlight", "우수", "best", "대표", "주요",
      "추천", "recommend"), "star"),
    (("흐름", "단계", "step", "flow", "순서", "다음", "진행", "next", "이동",
      "전환", "마이그레이션"), "chevron"),
    (("코드", "code", "개발", "develop", "빌드", "build", "git", "커밋",
      "commit", "푸시", "push"), "box"),
    (("사용자", "팀", "계정", "user", "team", "account", "people", "멤버",
      "member", "조직", "구성원"), "users"),
    (("완료", "확인", "승인", "check", "done", "complete", "approve", "검증",
      "verify", "성공"), "check"),
]

# 의미 키 → 네이티브 autoshape 멤버명(리치 아이콘 불가 시 폴백)
_KEY_TO_MSO = {
    "zap": "LIGHTNING_BOLT", "cloud": "CLOUD", "gear": "GEAR_6",
    "database": "CAN", "shield": "HEXAGON", "dollar": "DIAMOND",
    "chart": "DONUT", "file": "FOLDED_CORNER", "up": "UP_ARROW",
    "down": "DOWN_ARROW", "refresh": "CIRCULAR_ARROW", "star": "STAR_5_POINT",
    "chevron": "CHEVRON", "box": "CUBE", "users": "OVAL", "check": "OVAL",
    "server": "CUBE",
}


def _icon_name_for(text: str) -> str:
    """카드 텍스트 → 의미 아이콘 키. 미매칭 시 'chart'(파이)."""
    t = (text or "").lower()
    for keys, name in _ICON_MAP:
        for k in keys:
            if k in t:
                return name
    return "chart"


def build_native_diagram(slide, diagram_type: str, content: str,
                         region=(0.6, 1.7, 12.1, 5.2), palette=None,
                         title: str = "", backdrop: bool = False,
                         note: str = "") -> bool:
    """슬라이드에 네이티브(편집 가능) 다이어그램 도형을 그린다.

    Args:
        slide: python-pptx Slide
        diagram_type: "tree" | "flow" | "block" | (기타 → block)
        content: 다이어그램 원본 텍스트
        region: (left_in, top_in, width_in, height_in) 인치
        palette: ['#RRGGBB', ...] 선택적 색 팔레트(Style_Profile 유래)
        title: (옵션) 영역 상단 소제목

    Returns:
        True — 하나 이상의 도형을 그림 / False — 부재/빈콘텐츠/예외(호출자 PNG 폴백)
    """
    try:
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
        try:
            from pptx.enum.text import MSO_AUTO_SIZE
        except Exception:
            MSO_AUTO_SIZE = None
    except Exception:
        return False

    dtype = (diagram_type or "tree").lower()
    content = (content or "").strip()
    if not content:
        return False

    left_in, top_in, width_in, height_in = region
    # 제목이 있으면 상단에 공간을 확보(가상 y 92 위는 제목 영역).
    has_title = bool((title or "").strip())
    has_note = bool((note or "").strip())

    def _x(vx):
        return Inches(left_in + (vx / 100.0) * width_in)

    def _y(vy):
        return Inches(top_in + ((100.0 - vy) / 100.0) * height_in)

    def _w(vw):
        return Inches((vw / 100.0) * width_in)

    def _h(vh):
        return Inches((vh / 100.0) * height_in)

    pcols = _palette_colors(palette)
    primary = pcols[0] if len(pcols) >= 1 else "2E5BBA"
    secondary = pcols[1] if len(pcols) >= 2 else "5B9BD5"

    # 젠스파크풍 다채 셰이드(중복 제거, 순서 보존)
    _accent_shades = [primary, secondary, "2E5BBA", "3C78D8", "4A86E8",
                      "5B9BD5", "00897B", "F5A623", "7E57C2", "EF6C57"]
    _seen = set()
    _shades = []
    for _c in _accent_shades:
        cu = (_c or "").upper()
        if cu and cu not in _seen and _hex_to_rgb(cu):
            _seen.add(cu)
            _shades.append(cu)
    if not _shades:
        _shades = ["2E5BBA"]

    def _shade(i):
        return _shades[i % len(_shades)]

    def _rgb(hexstr):
        rgb = _hex_to_rgb(hexstr)
        return RGBColor(*rgb) if rgb else RGBColor(0x2E, 0x5B, 0xBA)

    def _tint(hexstr, factor=0.86):
        """배경(face)용 옅은 색."""
        rgb = _hex_to_rgb(hexstr) or (46, 91, 186)
        r, g, b = rgb
        return RGBColor(int(r + (255 - r) * factor),
                        int(g + (255 - g) * factor),
                        int(b + (255 - b) * factor))

    def _readable_text_rgb(hexstr):
        """배경 대비 가독 텍스트색 — 어두운 채움이면 흰색, 옅으면 짙은 회색."""
        rgb = _hex_to_rgb(hexstr) or (46, 91, 186)
        lum = 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]
        return RGBColor(0xFF, 0xFF, 0xFF) if lum < 140 else RGBColor(0x1A, 0x1A, 0x2A)

    def _shadow(shape):
        """부드러운 드롭 섀도우(outerShdw) XML 주입 — 깊이감(젠스파크풍)."""
        try:
            spPr = shape._element.spPr
            # 기존 effectLst 제거 후 재구성
            for el in spPr.findall(qn('a:effectLst')):
                spPr.remove(el)
            eff = spPr.makeelement(qn('a:effectLst'), {})
            sh = eff.makeelement(qn('a:outerShdw'), {
                'blurRad': '90000', 'dist': '38100', 'dir': '5400000',
                'rotWithShape': '0',
            })
            clr = sh.makeelement(qn('a:srgbClr'), {'val': '1A2A44'})
            alpha = clr.makeelement(qn('a:alpha'), {'val': '24000'})
            clr.append(alpha)
            sh.append(clr)
            eff.append(sh)
            spPr.append(eff)
        except Exception:
            pass

    def _set_text(shp, text, font_pt, bold, color_rgb, align=None, margin_left_in=None):
        tf = shp.text_frame
        tf.word_wrap = True
        if MSO_AUTO_SIZE is not None:
            try:
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            except Exception:
                pass
        try:
            # margin_left_in 지정 시 그만큼 좌측 여백 — 번호 배지 위에 텍스트가
            # 겹쳐 잘리는 것을 막는다(block/목차 가독성, 이슈2).
            tf.margin_left = Emu(int(margin_left_in * 914400)) if margin_left_in else Emu(54000)
            tf.margin_right = Emu(54000)
            tf.margin_top = Emu(27000)
            tf.margin_bottom = Emu(27000)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = align if align is not None else PP_ALIGN.CENTER
        run = p.add_run()
        run.text = _strip_md(str(text))[:220]
        run.font.size = Pt(font_pt)
        run.font.bold = bold
        run.font.color.rgb = color_rgb
        try:
            run.font.name = "Malgun Gothic"
        except Exception:
            pass

    def _card(vx, vy_top, vw, vh, text, fill_hex, edge_hex=None,
              font_pt=12, bold=False, solid_fill=False, shadow=True,
              text_rgb=None, align=None, radius=0.12, text_inset_left_in=None):
        """둥근 사각형 카드. solid_fill=True면 채움색=fill_hex(진한), 아니면 옅은 tint."""
        shp = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, _x(vx), _y(vy_top), _w(vw), _h(vh))
        # 모서리 반경 조절(adj)
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
        shp.fill.solid()
        if solid_fill:
            shp.fill.fore_color.rgb = _rgb(fill_hex)
            shp.line.color.rgb = _rgb(fill_hex)
            txt_rgb = text_rgb or _readable_text_rgb(fill_hex)
        else:
            shp.fill.fore_color.rgb = _tint(fill_hex)
            try:
                _spPr = shp._element.spPr
                for _e in _spPr.findall(qn('a:solidFill')):
                    _spPr.remove(_e)
                _rp = _hex_to_rgb(fill_hex) or (46, 91, 186)
                _fromhex = '%02X%02X%02X' % (int(_rp[0] + (255 - _rp[0]) * 0.80),
                                             int(_rp[1] + (255 - _rp[1]) * 0.80),
                                             int(_rp[2] + (255 - _rp[2]) * 0.80))
                _g = _spPr.makeelement(qn('a:gradFill'), {})
                _l = _g.makeelement(qn('a:gsLst'), {})
                for _pos, _col in ((0, _fromhex), (100000, 'FFFFFF')):
                    _gs = _l.makeelement(qn('a:gs'), {'pos': str(_pos)})
                    _cc = _gs.makeelement(qn('a:srgbClr'), {'val': _col})
                    _gs.append(_cc); _l.append(_gs)
                _g.append(_l)
                _g.append(_g.makeelement(qn('a:lin'), {'ang': '2700000', 'scaled': '1'}))
                _lnx = _spPr.find(qn('a:ln'))
                if _lnx is not None:
                    _lnx.addprevious(_g)
                else:
                    _spPr.append(_g)
            except Exception:
                pass
            shp.line.color.rgb = _rgb(edge_hex or fill_hex)
            txt_rgb = text_rgb or RGBColor(0x1A, 0x1A, 0x2A)
        shp.line.width = Pt(1.25)
        if shadow:
            _shadow(shp)
        _set_text(shp, text, font_pt, bold, txt_rgb, align=align,
                  margin_left_in=text_inset_left_in)
        return shp

    def _badge(cx, cy, d_v, label, fill_hex):
        """번호 원형 배지(편집 가능 oval)."""
        try:
            shp = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, _x(cx - d_v / 2.0), _y(cy + d_v / 2.0),
                _w(d_v), _h(d_v * (width_in / height_in)))
            shp.fill.solid()
            shp.fill.fore_color.rgb = _rgb(fill_hex)
            shp.line.fill.background()
            _set_text(shp, label, 11, True, RGBColor(0xFF, 0xFF, 0xFF))
            return shp
        except Exception:
            return None

    def _badge_in_gutter(card_vx, card_vy_top, card_vw, card_vh, d_v, label,
                         fill_hex, gutter="left"):
        """결함 C 수정 — 번호 배지를 라벨 카드 박스 '밖' 거터에 배치(순수 기하).

        spec: pptx-overlay-collision-fix, Task 3.3 / design Fix Implementation §2.

        라벨 카드 v-rect 를 인치로 환산(``_x``/``_y``/``_w``/``_h`` 역산)한 뒤
        ``layout_geometry.place_badge_in_gutter`` 로 카드 밖 거터의 배지 인치 좌표를
        구하고, 그 좌표를 ``_badge`` 의 중심 v-unit 으로 역산해 호출한다. 배지 사각형이
        라벨 카드 박스와 겹치지 않으므로 ``overlap_area(badge, label) == 0`` 이 보장된다.

        ``_badge`` 는 화면상 정사각 배지(지름 = ``(d_v/100)*width_in`` 인치)를 그리므로
        ``place_badge_in_gutter`` 에 같은 지름을 전달한다. 선호 거터(기본 ``"left"``)의
        결과가 슬라이드(13.333×7.5in) 밖으로 나가면 대체 거터로 폴백하고, 어떤 거터도
        담지 못하면 슬라이드 안으로 클램프한다(레이아웃 이탈 방지).
        """
        try:
            from .layout_geometry import place_badge_in_gutter
        except Exception:
            try:
                from ai_engine.layout_geometry import place_badge_in_gutter
            except Exception:
                from layout_geometry import place_badge_in_gutter

        # 라벨 카드 박스 인치 환산 (_x/_y/_w/_h 의 역산).
        cl = left_in + (card_vx / 100.0) * width_in
        ct = top_in + ((100.0 - card_vy_top) / 100.0) * height_in
        cw = (card_vw / 100.0) * width_in
        chh = (card_vh / 100.0) * height_in
        diameter_in = (d_v / 100.0) * width_in  # _badge 의 화면상 정사각 지름과 동일.
        card_rect = (cl, ct, cw, chh)

        _SW, _SH = 13.333, 7.5  # 표준 16:9 슬라이드 — 거터 결과 이탈 가드용.
        order = [gutter] + [g for g in ("left", "top", "bottom", "right")
                            if g != gutter]
        chosen = None
        for g in order:
            bl, bt, bw, bh = place_badge_in_gutter(
                card_rect, diameter_in, gutter=g, gap=0.05)
            if bl >= 0.0 and bt >= 0.0 and (bl + bw) <= _SW and (bt + bh) <= _SH:
                chosen = (bl, bt, bw, bh)
                break
        if chosen is None:
            # 모든 거터가 슬라이드 밖 → 선호 거터 결과를 슬라이드 안으로 클램프.
            bl, bt, bw, bh = place_badge_in_gutter(
                card_rect, diameter_in, gutter=gutter, gap=0.05)
            bl = min(max(bl, 0.0), max(0.0, _SW - bw))
            bt = min(max(bt, 0.0), max(0.0, _SH - bh))
            chosen = (bl, bt, bw, bh)
        bl, bt, bw, bh = chosen

        # _badge 중심 v-unit 역산: _x(cx - d_v/2) == bl, _y(cy + d_v/2) == bt.
        cx = d_v / 2.0 + (bl - left_in) * 100.0 / width_in
        cy = 100.0 - (bt - top_in) * 100.0 / height_in - d_v / 2.0
        return _badge(cx, cy, d_v, label, fill_hex)

    def _glyph(name, gx, gy_top, gw, gh, color=None):
        """의미 아이콘(흰색 솔리드 네이티브 도형)을 그린다. name은 MSO_SHAPE 멤버명."""
        try:
            enum = getattr(MSO_SHAPE, name, MSO_SHAPE.OVAL)
            g = slide.shapes.add_shape(enum, _x(gx), _y(gy_top), _w(gw), _h(gh))
            g.fill.solid()
            g.fill.fore_color.rgb = color or RGBColor(0xFF, 0xFF, 0xFF)
            g.line.fill.background()
            return g
        except Exception:
            return None

    def _connector(x1, y1, x2, y2, color_hex, arrow=True, width_pt=2.0):
        try:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, _x(x1), _y(y1), _x(x2), _y(y2))
            conn.line.color.rgb = _rgb(color_hex)
            conn.line.width = Pt(width_pt)
            if arrow:
                try:
                    ln = conn.line._get_or_add_ln()
                    tail = ln.makeelement(qn('a:tailEnd'),
                                          {'type': 'triangle', 'w': 'med', 'len': 'med'})
                    ln.append(tail)
                except Exception:
                    pass
            return conn
        except Exception:
            return None

    def _title_band():
        if not has_title:
            return
        # 업로드 PDF급 다크 섹션 헤더 밴드 — 좌측 컬러 strip + 차콜 바 + 흰 제목(편집 가능 도형).
        try:
            band = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, _x(2.0), _y(99.5), _w(96.0), _h(6.6))
            try:
                band.adjustments[0] = 0.14
            except Exception:
                pass
            band.fill.solid()
            band.fill.fore_color.rgb = RGBColor(0x2B, 0x38, 0x46)  # charcoal navy
            band.line.fill.background()
            try:
                _shadow(band)
            except Exception:
                pass
        except Exception:
            pass
        try:
            strip = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, _x(2.0), _y(99.5), _w(0.9), _h(6.6))
            strip.fill.solid()
            strip.fill.fore_color.rgb = _rgb(primary)
            strip.line.fill.background()
        except Exception:
            pass
        try:
            tb = slide.shapes.add_textbox(_x(4.2), _y(99.5), _w(92.5), _h(6.6))
            tf = tb.text_frame
            tf.word_wrap = True
            try:
                from pptx.enum.text import MSO_ANCHOR as _MA
                tf.vertical_anchor = _MA.MIDDLE
            except Exception:
                pass
            try:
                from pptx.util import Inches as _In2
                tf.margin_top = _In2(0.02)
                tf.margin_bottom = _In2(0.02)
            except Exception:
                pass
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(title)[:80]
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                run.font.name = "Malgun Gothic"
            except Exception:
                pass
        except Exception:
            pass

    # 제목이 있으면 다이어그램 본문을 아래로 눌러 겹침 방지.
    body_top = 90.0 if has_title else 100.0
    body_bot = 13.0 if has_note else 2.0

    # backdrop=True(슬라이드 배경 이미지 위에 그릴 때): 제목+다이어그램 영역을 덮는
    # 흰색 라운드 콘텐츠 카드를 먼저 깔아 가독성·편집성 확보(감마/젠스파크 카드 스타일).
    if backdrop:
        try:
            # 사방 여백을 둬 Vertex 배경이 카드를 프레임처럼 감싸 보이게 하고,
            # 카드는 반투명(프로스티드 글래스)으로 배경이 은은히 비치게 한다.
            # → 본문에서도 Vertex 고품질 배경이 분명히 드러난다(표지와 동일 품질감).
            try:
                _op = float(os.environ.get("AE_CARD_OPACITY", "0.86"))
            except Exception:
                _op = 0.86
            _op = max(0.55, min(0.97, _op))
            _alpha = str(int(_op * 100000))
            _SW, _SH = 13.333, 7.5
            margin = 0.5
            cx0 = margin
            ctop = 0.5
            card_w = _SW - 2 * margin
            ch = max(2.0, (_SH - 0.42) - ctop)
            panel = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx0), Inches(ctop),
                Inches(card_w), Inches(ch))
            try:
                panel.adjustments[0] = 0.045
            except Exception:
                pass
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            try:
                _sp = panel._element.spPr
                _sf = _sp.find(qn('a:solidFill'))
                _srgb = _sf.find(qn('a:srgbClr'))
                _al = _srgb.makeelement(qn('a:alpha'), {'val': _alpha})
                _srgb.append(_al)
            except Exception:
                pass
            panel.line.color.rgb = RGBColor(0xE2, 0xE6, 0xEE)
            panel.line.width = Pt(1.0)
            _shadow(panel)
        except Exception:
            pass

    drew = False
    try:
        if dtype == "flow":
            # content를 단계로 분해(라벨만 잘라내지 않고 설명까지 보존).
            _raw = content.replace("\u2192", "->").replace("\u21d2", "->")
            _steps_raw = []
            for _chunk in _raw.split("\n"):
                for _p in _chunk.split("->"):
                    _t = _p.strip().lstrip("-*\u2022").strip()
                    if _t:
                        _steps_raw.append(_t)
            _steps_raw = _steps_raw[:8]
            if not _steps_raw:
                return False

            def _split_step(s):
                """'제목: 설명' / '제목 - 설명' 등을 (제목, 설명)으로 분리."""
                s = str(s).strip()
                for sep in (" \u2014 ", " - ", " | ", ": ", "\uff1a", " :: ", "::", "\u2014", "|"):
                    if sep in s:
                        a, b = s.split(sep, 1)
                        a, b = a.strip(), b.strip()
                        if a and b:
                            return a[:40], b[:96]
                return s[:44], ""

            steps = [_split_step(s) for s in _steps_raw]
            _title_band()
            n = len(steps)
            _has_desc = any(d for _, d in steps)

            def _flow_card(vx, vy_top, vw, vh, ititle, idesc, sh, num, inset_in,
                           gutter="left"):
                """단계 카드 — 제목(굵게) + 설명(작게, 회색)을 함께 그린다."""
                box = _card(vx, vy_top, vw, vh, ititle, sh, font_pt=13, bold=True,
                            solid_fill=False, align=PP_ALIGN.LEFT, radius=0.16,
                            text_inset_left_in=inset_in)
                if idesc:
                    try:
                        tf = box.text_frame
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        p = tf.add_paragraph()
                        p.alignment = PP_ALIGN.LEFT
                        r = p.add_run()
                        r.text = idesc
                        r.font.size = Pt(10)
                        r.font.bold = False
                        r.font.color.rgb = RGBColor(0x5A, 0x60, 0x6A)
                        try:
                            r.font.name = "Malgun Gothic"
                        except Exception:
                            pass
                    except Exception:
                        pass
                # 결함 C 수정 — 번호 배지를 라벨 카드 밖 거터로(카드 v-rect 기준).
                _badge_in_gutter(vx, vy_top, vw, vh, 5.4, str(num), sh, gutter=gutter)
                return box

            if n <= 4:
                # ── 가로 흐름 — 균등 폭 카드 + 번호 배지 + 제목/설명 + 정밀 화살표 ──
                margin = 3.0
                arrow_gap = 4.0
                usable = 100.0 - 2 * margin - (n - 1) * arrow_gap
                w = usable / n
                box_h = 40.0 if _has_desc else 26.0
                cy_top = (body_top + body_bot) / 2.0 + box_h / 2.0
                ymid = cy_top - box_h / 2.0
                _inset = max(0.0, (5.4 / 100.0) * width_in + 0.18)
                x = margin
                centers = []
                for i, (ititle, idesc) in enumerate(steps):
                    sh = _shade(i)
                    _flow_card(x, cy_top, w, box_h, ititle, idesc, sh, i + 1, _inset,
                               gutter="top")
                    centers.append((x, x + w))
                    drew = True
                    x += w + arrow_gap
                for i in range(n - 1):
                    _connector(centers[i][1], ymid, centers[i + 1][0], ymid,
                               _shade(i), arrow=True, width_pt=2.25)
            else:
                # ── 세로 흐름(5+단계) — 좌측 번호 배지 + 넓은 카드(제목/설명) + 아래 화살표 ──
                top_v, bot_v = body_top - 2, body_bot
                slot = (top_v - bot_v) / n
                box_h = max(8.0, min(slot - 2.4, 15.0))
                box_w = 82.0
                x = 9.0
                _inset = max(0.0, (5.4 / 100.0) * box_w / 100.0 * width_in + 0.55)
                col = []
                for i, (ititle, idesc) in enumerate(steps):
                    vy_t = top_v - i * slot
                    sh = _shade(i)
                    _flow_card(x, vy_t, box_w, box_h, ititle, idesc, sh, i + 1, _inset,
                               gutter="left")
                    col.append((x + box_w / 2.0, vy_t - box_h))
                    drew = True
                for i in range(n - 1):
                    cx = col[i][0]
                    _connector(cx, col[i][1], cx, top_v - (i + 1) * slot,
                               _shade(i), arrow=True, width_pt=2.25)

        elif dtype == "tree":
            nodes = _parse_tree(content)
            if not nodes:
                return False
            _title_band()
            roots = [nd for nd in nodes if nd["depth"] == 0]
            children0 = [nd for nd in nodes if nd["depth"] == 1]
            deeper = [nd for nd in nodes if nd["depth"] >= 2]

            if len(roots) == 1 and children0 and not deeper and 2 <= len(children0) <= 8:
                # ── 조직도(org-chart) — 루트 상단 중앙(진한 채움) + 자식 균등 + 엘보 ──
                root = roots[0]
                nc = len(children0)
                root_w = max(22.0, min(40.0, _text_units(root["text"]) * 2.4 + 10))
                root_h = 15.0
                root_top = body_top
                root_cx = 50.0
                _card(root_cx - root_w / 2.0, root_top, root_w, root_h,
                      root["text"], primary, font_pt=15, bold=True,
                      solid_fill=True, radius=0.16)
                drew = True
                root_bottom = root_top - root_h

                margin = 2.0
                gap = 2.6
                child_w = min(22.0, (100.0 - 2 * margin - (nc - 1) * gap) / nc)
                total_w = nc * child_w + (nc - 1) * gap
                start_x = (100.0 - total_w) / 2.0
                child_h = 26.0
                child_top = root_bottom - 16.0
                bus_y = (root_bottom + child_top) / 2.0
                centers = []
                for i, ch in enumerate(children0):
                    cx_left = start_x + i * (child_w + gap)
                    ccx = cx_left + child_w / 2.0
                    centers.append(ccx)
                    sh = _shade(i)
                    _card(cx_left, child_top, child_w, child_h, ch["text"], sh,
                          font_pt=12, bold=False, solid_fill=False, radius=0.16)
                    _connector(ccx, bus_y, ccx, child_top, sh, arrow=True, width_pt=2.0)
                _connector(root_cx, root_bottom, root_cx, bus_y, primary,
                           arrow=False, width_pt=2.0)
                xs = centers + [root_cx]
                _connector(min(xs), bus_y, max(xs), bus_y, primary,
                           arrow=False, width_pt=2.0)
                return drew

            # ── 세로 들여쓰기 트리(파일/깊은 구조) ──
            n = len(nodes)
            top, bottom = body_top - 2, body_bot
            avail = top - bottom
            box_h = max(5.0, min(9.0, avail / max(n, 1) - 1.2))
            y_step = avail / max(n, 1)
            indent_unit = 6.5
            base_x = 2.0
            positions = []
            for i, nd in enumerate(nodes):
                vx = base_x + nd["depth"] * indent_unit
                vy_t = top - i * y_step
                positions.append({
                    "x": vx, "y_top": vy_t, "y_bot": vy_t - box_h,
                    "y_center": vy_t - box_h / 2,
                    "depth": nd["depth"], "is_dir": nd["is_dir"], "text": nd["text"],
                })
            for i, pp in enumerate(positions):
                if pp["depth"] == 0:
                    continue
                parent = None
                for j in range(i - 1, -1, -1):
                    if positions[j]["depth"] < pp["depth"]:
                        parent = positions[j]
                        break
                if not parent:
                    continue
                x_branch = parent["x"] + 1.8
                _connector(x_branch, parent["y_bot"], x_branch, pp["y_center"],
                           secondary, arrow=False, width_pt=1.5)
                _connector(x_branch, pp["y_center"], pp["x"], pp["y_center"],
                           secondary, arrow=False, width_pt=1.5)
            for pp in positions:
                if pp["is_dir"]:
                    label = pp["text"][:46]
                    sh = primary
                    bold = True
                else:
                    label = pp["text"][:48]
                    sh = secondary
                    bold = False
                w_eff = max(20.0, 36.0 - pp["depth"] * 2.5)
                _card(pp["x"], pp["y_top"], w_eff, box_h, label, sh,
                      font_pt=10, bold=bold, solid_fill=False, shadow=False,
                      align=PP_ALIGN.LEFT, radius=0.2)
                drew = True

        elif dtype == "cards":
            # ── 젠스파크급 피처 카드 그리드 (편집 가능 네이티브) ──
            # 그라데이션 카드 + 좌측 컬러 액센트 바 + 솔리드 아이콘 칩(흰 마크) + 굵은 제목 + 옅은 설명.
            cards = _parse_cards(content)
            if not cards:
                return False
            _title_band()
            n = len(cards)
            cols = 3 if n >= 5 else (2 if n >= 2 else 1)
            rows = (n + cols - 1) // cols
            m_x = 3.0
            gap_x, gap_y = 3.0, 4.0
            top_v = body_top - 1.0
            avail_h = top_v - body_bot
            card_w = (100.0 - 2 * m_x - (cols - 1) * gap_x) / cols
            card_h = min(36.0, (avail_h - (rows - 1) * gap_y) / max(rows, 1))
            asp = width_in / height_in  # 1 virtual-x : asp virtual-y = 화면 정사각

            def _hexstr(hexstr, factor):
                rgb = _hex_to_rgb(hexstr) or (46, 91, 186)
                r, g, b = rgb
                return "%02X%02X%02X" % (int(r + (255 - r) * factor),
                                          int(g + (255 - g) * factor),
                                          int(b + (255 - b) * factor))

            def _grad_fill(shp, hex_from):
                """좌상→우하 옅은 그라데이션 배경 주입(편집 가능 도형 유지)."""
                try:
                    spPr = shp._element.spPr
                    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
                        for el in spPr.findall(qn(tag)):
                            spPr.remove(el)
                    grad = spPr.makeelement(qn("a:gradFill"), {})
                    lst = grad.makeelement(qn("a:gsLst"), {})
                    for pos, col in ((0, hex_from), (100000, "FFFFFF")):
                        gs = lst.makeelement(qn("a:gs"), {"pos": str(pos)})
                        cc = gs.makeelement(qn("a:srgbClr"), {"val": col})
                        gs.append(cc)
                        lst.append(gs)
                    grad.append(lst)
                    grad.append(grad.makeelement(qn("a:lin"), {"ang": "2700000", "scaled": "1"}))
                    ln = spPr.find(qn("a:ln"))
                    if ln is not None:
                        ln.addprevious(grad)
                    else:
                        spPr.append(grad)
                    return True
                except Exception:
                    try:
                        shp.fill.solid()
                        shp.fill.fore_color.rgb = _tint(hex_from, 0.88)
                    except Exception:
                        pass
                    return False

            for idx, (ctitle, cdesc) in enumerate(cards):
                r = idx // cols
                c = idx % cols
                cx = m_x + c * (card_w + gap_x)
                cy_top = top_v - r * (card_h + gap_y)
                sh = _shade(idx)
                # 카드 본체 — 옅은 그라데이션 + 부드러운 섀도우(테두리 없음)
                cardshp = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, _x(cx), _y(cy_top), _w(card_w), _h(card_h))
                try:
                    cardshp.adjustments[0] = 0.07
                except Exception:
                    pass
                _grad_fill(cardshp, _hexstr(sh, 0.90))
                try:
                    cardshp.line.color.rgb = _rgb(_hexstr(sh, 0.62))
                    cardshp.line.width = Pt(1.0)
                except Exception:
                    pass
                _shadow(cardshp)
                # 좌측 컬러 액센트 바
                try:
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, _x(cx), _y(cy_top), _w(0.8), _h(card_h))
                    bar.adjustments[0] = 0.5
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = _rgb(sh)
                    bar.line.fill.background()
                except Exception:
                    pass
                pad = card_w * 0.10
                # 아이콘 칩 — 솔리드 컬러 둥근 사각형 + 흰 원 마크 (화면상 정사각)
                chip_w = min(card_w * 0.24, 7.5)
                chip_h = chip_w * asp
                chip_top = cy_top - card_h * 0.10
                try:
                    chipshp = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        _x(cx + pad), _y(chip_top), _w(chip_w), _h(chip_h))
                    chipshp.adjustments[0] = 0.30
                    chipshp.fill.solid()
                    chipshp.fill.fore_color.rgb = _rgb(sh)
                    chipshp.line.fill.background()
                    _shadow(chipshp)
                    # 의미 아이콘 — 1순위 Lucide 고품질 벡터(흰색 PNG), 폴백 네이티브 autoshape.
                    _ikey = _icon_name_for((ctitle or "") + " " + (cdesc or ""))
                    mk = chip_w * 0.54
                    _ic_x = cx + pad + (chip_w - mk) / 2
                    _ic_top = chip_top - (chip_h - mk * asp) / 2
                    _ic_png = None
                    try:
                        import os as _os_ic
                        if _os_ic.environ.get("AE_DISABLE_RICH_ICONS", "") != "1":
                            try:
                                from icon_assets import get_icon_png as _gip
                            except Exception:
                                from ai_engine.icon_assets import get_icon_png as _gip
                            _ic_png = _gip(_ikey, "FFFFFF")
                    except Exception:
                        _ic_png = None
                    # D2(슬롯-이미지 크기 정합): 아이콘 칩은 소형 장식 슬롯이다.
                    # 대형 이미지가 자산 경로로 흔러들면 짤그러지므로, 임베드 직전 픽셀
                    # 크기를 PIL로 측정해 slot_image_fits 로 검사한다. 소형 슬롯+대형 이미지면
                    # PNG 를 임베드하지 않고 네이티브 글리프로 폴백한다(소형 슬롯에 대형 이미지 방지).
                    _ic_fits = True
                    if _ic_png:
                        try:
                            from PIL import Image as _PILic
                            with _PILic.open(_ic_png) as _imic:
                                _icpw, _icph = _imic.size
                            try:
                                from .layout_geometry import slot_image_fits as _sif_ic
                            except Exception:
                                try:
                                    from ai_engine.layout_geometry import slot_image_fits as _sif_ic
                                except Exception:
                                    from layout_geometry import slot_image_fits as _sif_ic
                            _slot_w_in = (mk / 100.0) * width_in
                            _slot_h_in = (mk * asp / 100.0) * height_in
                            _ic_fits = _sif_ic((0.0, 0.0, _slot_w_in, _slot_h_in), _icpw, _icph)
                        except Exception:
                            _ic_fits = True
                    if _ic_png and _ic_fits:
                        try:
                            slide.shapes.add_picture(
                                _ic_png, _x(_ic_x), _y(_ic_top), _w(mk), _h(mk * asp))
                        except Exception:
                            _glyph(_KEY_TO_MSO.get(_ikey, "DONUT"), _ic_x, _ic_top, mk, mk * asp)
                    else:
                        _glyph(_KEY_TO_MSO.get(_ikey, "DONUT"), _ic_x, _ic_top, mk, mk * asp)
                except Exception:
                    pass
                # 제목(굵게) — 칩 아래
                title_top = chip_top - chip_h - card_h * 0.04
                try:
                    tt = slide.shapes.add_textbox(
                        _x(cx + pad), _y(title_top), _w(card_w - 2 * pad), _h(card_h * 0.22))
                    _set_text(tt, ctitle, 15, True, RGBColor(0x1A, 0x1A, 0x2A), align=PP_ALIGN.LEFT)
                except Exception:
                    pass
                # 설명(옅은 회색) — 제목 아래
                if cdesc:
                    try:
                        dd = slide.shapes.add_textbox(
                            _x(cx + pad), _y(title_top - card_h * 0.24),
                            _w(card_w - 2 * pad), _h(card_h * 0.36))
                        _set_text(dd, cdesc, 11, False, RGBColor(0x6B, 0x72, 0x80), align=PP_ALIGN.LEFT)
                    except Exception:
                        pass
                drew = True

        elif dtype == "twocol":
            # ── 2단 컬럼 — 항목을 좌/우 두 열에 번호 배지 + 제목(+설명) 카드로 쌓는다.
            #    doc2(IT 온보딩)의 2단 밀도 재현. 각 카드 편집 가능. ──
            items = _parse_cards(content)
            if not items:
                return False
            _title_band()
            n = len(items)
            half = (n + 1) // 2
            cols_data = [items[:half], items[half:]]
            col_w = 45.0
            gap = 6.0
            xs = [2.5, 2.5 + 45.0 + 6.0]
            top_v = body_top - 1.0
            avail_h = top_v - body_bot
            base_idx = 0
            for _ci2, col_items in enumerate(cols_data):
                if not col_items:
                    continue
                rows = len(col_items)
                ch = min(22.0, (avail_h - (rows - 1) * 3.0) / max(rows, 1))
                x0 = xs[_ci2]
                for r, (ct, cd) in enumerate(col_items):
                    vy_t = top_v - r * (ch + 3.0)
                    sh = _shade(base_idx + r)
                    _shp2 = _card(x0, vy_t, col_w, ch, ct, sh, font_pt=13, bold=True,
                                  solid_fill=False, align=PP_ALIGN.LEFT, radius=0.12,
                                  text_inset_left_in=0.95)
                    if cd and _shp2 is not None:
                        try:
                            _tf2 = _shp2.text_frame
                            _tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
                            _pp2 = _tf2.add_paragraph(); _pp2.alignment = PP_ALIGN.LEFT
                            _rr2 = _pp2.add_run(); _rr2.text = _strip_md(cd)[:90]
                            _rr2.font.size = Pt(10); _rr2.font.bold = False
                            _rr2.font.color.rgb = RGBColor(0x5A, 0x60, 0x6A)
                            try:
                                _rr2.font.name = "Malgun Gothic"
                            except Exception:
                                pass
                        except Exception:
                            pass
                    # 좌측 컬러 액센트 바
                    try:
                        bar2 = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE, _x(x0), _y(vy_t), _w(0.8), _h(ch))
                        bar2.adjustments[0] = 0.5
                        bar2.fill.solid(); bar2.fill.fore_color.rgb = _rgb(sh)
                        bar2.line.fill.background()
                    except Exception:
                        pass
                    # 결함 C 수정 — 번호 배지를 라벨 카드 밖 좌측 거터로.
                    _badge_in_gutter(x0, vy_t, col_w, ch, 5.6,
                                     str(base_idx + r + 1), sh, gutter="left")
                base_idx += rows
            drew = True

        elif dtype in ("architecture", "stack"):
            # ── 계층형 아키텍처/스택 — 중앙 정렬 풀폭 타일을 위→아래로 쌓고
            #    얇은 스파인 커넥터로 연결. block(좌측 배지 리스트)보다 '계층/스택'
            #    의미를 더 잘 전달한다. 각 타일 편집 가능. ──
            lines = _parse_block(content)
            if not lines:
                return False
            _title_band()
            n = len(lines)
            top_v, bot_v = body_top, body_bot
            slot = (top_v - bot_v) / n
            band_h = max(7.0, min(slot - 3.0, 15.0))
            band_w = 80.0
            x = (100.0 - band_w) / 2.0
            spans = []
            for i, txt in enumerate(lines):
                vy_t = top_v - i * slot
                sh = _shade(i)
                # tint 배경 + 색 테두리 + 중앙 굵은 텍스트(계층 라벨)
                _card(x, vy_t, band_w, band_h, txt, sh, font_pt=14, bold=True,
                      solid_fill=False, align=PP_ALIGN.CENTER, radius=0.10)
                # 좌측 색 스트라이프(계층 강조, 편집 가능)
                try:
                    stripe = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, _x(x), _y(vy_t),
                        _w(1.6), _h(band_h))
                    stripe.fill.solid()
                    stripe.fill.fore_color.rgb = _rgb(sh)
                    stripe.line.fill.background()
                except Exception:
                    pass
                spans.append((x + band_w / 2.0, vy_t, vy_t - band_h))
                drew = True
            # 타일 사이 얇은 스파인 커넥터(계층 흐름)
            for i in range(n - 1):
                cx = spans[i][0]
                _connector(cx, spans[i][2], cx, spans[i + 1][1], _shade(i),
                           arrow=False, width_pt=1.75)

        elif dtype == "kpi":
            # ── KPI/지표 요약 카드 (젠스파크 OKR 스타일, 편집 가능) ──
            # 큰 숫자 + 라벨 + (옵션)증감 배지를 가로로 나열.
            kpis = _parse_kpis(content)
            if not kpis:
                return False
            _title_band()
            cols = min(len(kpis), 4)
            m_x, gap_x = 3.0, 3.0
            card_w = (100.0 - 2 * m_x - (cols - 1) * gap_x) / cols
            card_h = min(34.0, (body_top - body_bot) * 0.66)
            cy_top = body_top - (body_top - body_bot - card_h) / 2.0
            for idx, (value, label, delta) in enumerate(kpis[:cols]):
                cx = m_x + idx * (card_w + gap_x)
                sh = _shade(idx)
                card = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, _x(cx), _y(cy_top), _w(card_w), _h(card_h))
                try:
                    card.adjustments[0] = 0.08
                except Exception:
                    pass
                card.fill.solid()
                card.fill.fore_color.rgb = _tint(sh, 0.92)
                card.line.color.rgb = _rgb(sh)
                card.line.width = Pt(1.5)
                _shadow(card)
                # 상단 색 띠(액센트)
                try:
                    cap = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, _x(cx), _y(cy_top), _w(card_w), _h(2.4))
                    cap.adjustments[0] = 0.5
                    cap.fill.solid()
                    cap.fill.fore_color.rgb = _rgb(sh)
                    cap.line.fill.background()
                except Exception:
                    pass
                # 큰 값
                vbox = slide.shapes.add_textbox(
                    _x(cx), _y(cy_top - card_h * 0.12), _w(card_w), _h(card_h * 0.46))
                _set_text(vbox, value or label, 34, True, _rgb(sh), align=PP_ALIGN.CENTER)
                # 라벨
                if value:
                    lbox = slide.shapes.add_textbox(
                        _x(cx + card_w * 0.04), _y(cy_top - card_h * 0.56),
                        _w(card_w * 0.92), _h(card_h * 0.3))
                    _set_text(lbox, label, 12, False, RGBColor(0x44, 0x4A, 0x55), align=PP_ALIGN.CENTER)
                # 증감 배지
                if delta:
                    up = not delta.lstrip("+").startswith(("-", "\u25bc", "\u25bd"))
                    dcol = "1A8A4A" if up else "C0392B"
                    try:
                        db = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            _x(cx + card_w * 0.26), _y(cy_top - card_h * 0.84),
                            _w(card_w * 0.48), _h(card_h * 0.15))
                        db.adjustments[0] = 0.5
                        db.fill.solid()
                        db.fill.fore_color.rgb = _tint(dcol, 0.82)
                        db.line.color.rgb = _rgb(dcol)
                        db.line.width = Pt(1.0)
                        _set_text(db, delta, 10, True, _rgb(dcol), align=PP_ALIGN.CENTER)
                    except Exception:
                        pass
                drew = True

        elif dtype == "progress":
            # ── 진행률/목표 달성 막대 (젠스파크 상태 테이블 스타일, 편집 가능) ──
            # 각 행: 라벨(좌) + 진행 막대(트랙+채움) + % + (옵션)상태 칩.
            rows_data = _parse_progress(content)
            if not rows_data:
                return False
            _title_band()
            n = len(rows_data)
            top_v, bot_v = body_top, body_bot
            slot = (top_v - bot_v) / n
            row_h = max(6.0, min(slot - 2.2, 13.0))
            label_w, pct_w, status_w = 30.0, 8.0, 14.0
            bar_x = label_w + 2.0
            bar_w = 100.0 - label_w - pct_w - status_w - 6.0
            track_h = max(3.0, row_h * 0.40)
            for i, (label, pct, status) in enumerate(rows_data):
                vy_t = top_v - i * slot
                sh = _shade(i)
                cy_mid = vy_t - row_h / 2.0
                # 라벨(좌)
                lb = slide.shapes.add_textbox(_x(1.0), _y(vy_t), _w(label_w), _h(row_h))
                _set_text(lb, label, 13, True, RGBColor(0x1A, 0x1A, 0x2A), align=PP_ALIGN.LEFT)
                # 트랙
                try:
                    track = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        _x(bar_x), _y(cy_mid + track_h / 2.0), _w(bar_w), _h(track_h))
                    track.adjustments[0] = 0.5
                    track.fill.solid()
                    track.fill.fore_color.rgb = RGBColor(0xE6, 0xE9, 0xEF)
                    track.line.fill.background()
                except Exception:
                    pass
                # 채움(pct 비율)
                fill_w = max(1.5, bar_w * pct / 100.0)
                try:
                    fillbar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        _x(bar_x), _y(cy_mid + track_h / 2.0), _w(fill_w), _h(track_h))
                    fillbar.adjustments[0] = 0.5
                    fillbar.fill.solid()
                    fillbar.fill.fore_color.rgb = _rgb(sh)
                    fillbar.line.fill.background()
                    _shadow(fillbar)
                except Exception:
                    pass
                # % 텍스트
                pb = slide.shapes.add_textbox(
                    _x(bar_x + bar_w + 1.0), _y(vy_t), _w(pct_w + 3.0), _h(row_h))
                _set_text(pb, f"{pct}%", 13, True, _rgb(sh), align=PP_ALIGN.LEFT)
                # 상태 칩(옵션)
                if status:
                    try:
                        chip = slide.shapes.add_shape(
                            MSO_SHAPE.ROUNDED_RECTANGLE,
                            _x(100.0 - status_w), _y(cy_mid + row_h * 0.30),
                            _w(status_w - 1.0), _h(row_h * 0.6))
                        chip.adjustments[0] = 0.5
                        chip.fill.solid()
                        chip.fill.fore_color.rgb = _tint(sh, 0.80)
                        chip.line.color.rgb = _rgb(sh)
                        chip.line.width = Pt(1.0)
                        _set_text(chip, status, 10, True, _rgb(sh), align=PP_ALIGN.CENTER)
                    except Exception:
                        pass
                drew = True

        else:  # block — 번호 배지 + 좌측 액센트 바 + 카드(균등 세로 배치)
            lines = _parse_block(content)
            if not lines:
                return False
            _title_band()
            n = len(lines)
            top_v, bot_v = body_top, body_bot
            slot = (top_v - bot_v) / n
            box_h = max(8.0, min(slot - 4.0, 16.0))
            box_w = 86.0
            x = 7.0

            def _split_bl(s):
                s = str(s).strip()
                for sep in (" \u2014 ", " - ", " | ", ": ", "\uff1a", "\u2014", "|"):
                    if sep in s:
                        a, b = s.split(sep, 1)
                        a, b = a.strip(), b.strip()
                        if a and b:
                            return a[:50], b[:110]
                return s[:70], ""

            for i, txt in enumerate(lines):
                vy_t = top_v - i * slot
                sh = _shade(i)
                # 카드 — 좌측 번호 배지/액센트 바를 피해 텍스트를 들여쓴다.
                # '제목: 설명'이면 제목(굵게)+설명(작게/회색)으로 분리 표시(레퍼런스 스타일).
                _bt, _bd = _split_bl(txt)
                _bcard = _card(x, vy_t, box_w, box_h, _bt, sh, font_pt=13,
                      bold=True, solid_fill=False, align=PP_ALIGN.LEFT, radius=0.12,
                      text_inset_left_in=0.95)
                if _bd:
                    try:
                        _btf = _bcard.text_frame
                        _btf.vertical_anchor = MSO_ANCHOR.MIDDLE
                        _bp = _btf.add_paragraph()
                        _bp.alignment = PP_ALIGN.LEFT
                        _br = _bp.add_run()
                        _br.text = _strip_md(_bd)
                        _br.font.size = Pt(10)
                        _br.font.bold = False
                        _br.font.color.rgb = RGBColor(0x5A, 0x60, 0x6A)
                        try:
                            _br.font.name = "Malgun Gothic"
                        except Exception:
                            pass
                    except Exception:
                        pass
                # 좌측 액센트 바(편집 가능 사각형)
                try:
                    bar = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE, _x(x), _y(vy_t),
                        _w(1.4), _h(box_h))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = _rgb(sh)
                    bar.line.fill.background()
                except Exception:
                    pass
                # 결함 C 수정 — 번호 배지를 라벨 카드 밖 좌측 거터로.
                _badge_in_gutter(x, vy_t, box_w, box_h, 6.2, str(i + 1), sh,
                                 gutter="left")
                drew = True
    except Exception:
        return drew

    # NOTICE 콜아웃 박스 (크림 배경 + 주황 좌측 보더) — doc2 시그니처. 편집 가능.
    if has_note:
        try:
            _cn_top, _cn_h = 11.5, 9.0
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, _x(2.0), _y(_cn_top), _w(96.0), _h(_cn_h))
            try:
                box.adjustments[0] = 0.08
            except Exception:
                pass
            box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xFF, 0xF8, 0xEE)
            box.line.color.rgb = RGBColor(0xF0, 0xD8, 0xB8); box.line.width = Pt(1.0)
            try:
                _shadow(box)
            except Exception:
                pass
            lb = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, _x(2.0), _y(_cn_top), _w(0.7), _h(_cn_h))
            lb.fill.solid(); lb.fill.fore_color.rgb = RGBColor(0xED, 0x7D, 0x31)
            lb.line.fill.background()
            tb = slide.shapes.add_textbox(_x(3.6), _y(_cn_top), _w(92.0), _h(_cn_h))
            tf = tb.text_frame; tf.word_wrap = True
            try:
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            except Exception:
                pass
            p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.LEFT
            r0 = p0.add_run(); r0.text = "NOTICE"
            r0.font.size = Pt(11); r0.font.bold = True
            r0.font.color.rgb = RGBColor(0xED, 0x7D, 0x31)
            try:
                r0.font.name = "Malgun Gothic"
            except Exception:
                pass
            p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.LEFT
            r1 = p1.add_run(); r1.text = _strip_md(str(note))[:200]
            r1.font.size = Pt(13); r1.font.bold = False
            r1.font.color.rgb = RGBColor(0x5A, 0x4A, 0x33)
            try:
                r1.font.name = "Malgun Gothic"
            except Exception:
                pass
            drew = True
        except Exception:
            pass

    return drew


def build_native_cover(slide, prs, title: str, subtitle: str = "",
                       eyebrow: str = "", date_str: str = "",
                       kpis=None, palette=None,
                       remove_placeholders: bool = True,
                       over_image: bool = False) -> bool:
    """표지(cover) 슬라이드를 편집 가능한 네이티브 도형으로 디자인한다.

    템플릿 표지는 제목 placeholder만 채워 거의 비어 보이는(채움 ~3%) 문제가 있다.
    이 함수는 좌측 액센트 바 + 아이브로우 룰 + 대형 제목 + 부제 + 날짜/지표 +
    우상단 기하 액센트 클러스터로 젠스파크 OKR 표지급 비주얼을 네이티브로 조립한다.

    모든 요소는 PowerPoint에서 개별 편집 가능(통짜 PNG 아님).

    Args:
        slide: python-pptx Slide (표지)
        prs: Presentation (슬라이드 치수 산출용)
        title: 표지 대제목
        subtitle: 부제(옵션)
        eyebrow: 상단 소형 라벨(옵션, 대문자 권장)
        date_str: 날짜 문자열(옵션)
        kpis: [(value, label, delta), ...] 하단 KPI 카드(옵션, 최대 4)
        palette: ['#RRGGBB', ...] 팔레트(첫 색=primary)
        remove_placeholders: True면 기존 title/date placeholder 제거(중복/겹침 방지)

    Returns:
        True — 하나 이상 도형을 그림 / False — 예외(호출자 폴백)
    """
    try:
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.oxml.ns import qn
    except Exception:
        return False

    import math
    # 결함 A 수정 — 제목/부제 비겹침 수직 스택(순수 기하). import 컨텍스트에 견고하게.
    try:
        from .layout_geometry import vertical_stack
    except Exception:
        try:
            from ai_engine.layout_geometry import vertical_stack
        except Exception:
            from layout_geometry import vertical_stack

    title = (title or "").strip()
    if not title:
        return False

    # 슬라이드 치수(인치)
    try:
        SW = prs.slide_width / 914400.0
        SH = prs.slide_height / 914400.0
    except Exception:
        SW, SH = 13.333, 7.5

    pcols = _palette_colors(palette)
    primary = pcols[0] if len(pcols) >= 1 else "2E5BBA"
    secondary = pcols[1] if len(pcols) >= 2 else "5B9BD5"
    accent = pcols[2] if len(pcols) >= 3 else secondary

    def _rgb(hexstr):
        rgb = _hex_to_rgb(hexstr)
        return RGBColor(*rgb) if rgb else RGBColor(0x2E, 0x5B, 0xBA)

    def _tint(hexstr, factor=0.86):
        rgb = _hex_to_rgb(hexstr) or (46, 91, 186)
        r, g, b = rgb
        return RGBColor(int(r + (255 - r) * factor),
                        int(g + (255 - g) * factor),
                        int(b + (255 - b) * factor))

    def _readable(hexstr):
        rgb = _hex_to_rgb(hexstr) or (46, 91, 186)
        lum = 0.299 * rgb[0] + 0.587 * rgb[1] + 0.114 * rgb[2]
        return RGBColor(0xFF, 0xFF, 0xFF) if lum < 150 else RGBColor(0x1A, 0x1A, 0x2A)

    def _shadow(shape, alpha="22000", blur="80000"):
        try:
            spPr = shape._element.spPr
            for el in spPr.findall(qn('a:effectLst')):
                spPr.remove(el)
            eff = spPr.makeelement(qn('a:effectLst'), {})
            sh = eff.makeelement(qn('a:outerShdw'), {
                'blurRad': blur, 'dist': '34000', 'dir': '5400000', 'rotWithShape': '0'})
            clr = sh.makeelement(qn('a:srgbClr'), {'val': '1A2A44'})
            a = clr.makeelement(qn('a:alpha'), {'val': alpha})
            clr.append(a); sh.append(clr); eff.append(sh); spPr.append(eff)
        except Exception:
            pass

    def _grad(shp, hex_from, hex_to, ang="2700000"):
        try:
            spPr = shp._element.spPr
            for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:blipFill", "a:pattFill"):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
            g = spPr.makeelement(qn("a:gradFill"), {})
            lst = g.makeelement(qn("a:gsLst"), {})
            for pos, col in ((0, hex_from), (100000, hex_to)):
                gs = lst.makeelement(qn("a:gs"), {"pos": str(pos)})
                cc = gs.makeelement(qn("a:srgbClr"), {"val": col})
                gs.append(cc); lst.append(gs)
            g.append(lst)
            g.append(g.makeelement(qn("a:lin"), {"ang": ang, "scaled": "1"}))
            ln = spPr.find(qn("a:ln"))
            if ln is not None:
                ln.addprevious(g)
            else:
                spPr.append(g)
        except Exception:
            pass

    def _txt(x, y, w, h, text, pt, bold, color, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, spacing=None):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        try:
            tf.vertical_anchor = anchor
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        except Exception:
            pass
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(pt)
        run.font.bold = bold
        run.font.color.rgb = color
        try:
            run.font.name = "Malgun Gothic"
        except Exception:
            pass
        if spacing is not None:
            try:
                from pptx.oxml.ns import qn as _qn
                rPr = run._r.get_or_add_rPr()
                rPr.set('spc', str(int(spacing)))
            except Exception:
                pass
        return tb

    # 0) 기존 placeholder + 텍스트 보유 도형 제거(제목 중복/겹침 방지).
    #    그림(템플릿 배경·로고)과 빈 장식 도형은 보존한다.
    if remove_placeholders:
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
        except Exception:
            _MST = None
        try:
            for shp in list(slide.shapes):
                is_pic = bool(_MST is not None and shp.shape_type == _MST.PICTURE)
                try:
                    is_ph = bool(shp.is_placeholder)
                except Exception:
                    is_ph = False
                try:
                    has_text = bool(shp.has_text_frame and shp.text_frame.text.strip())
                except Exception:
                    has_text = False
                if is_ph or (has_text and not is_pic):
                    try:
                        shp._element.getparent().remove(shp._element)
                    except Exception:
                        pass
        except Exception:
            pass

    # 0.5) 이미지 위에 그릴 때(over_image): 어두운 스크림을 깔아 흰 텍스트 가독성 확보.
    #      좌측이 더 어두운 그라데이션 → 우측 이미지가 더 보이게.
    if over_image:
        try:
            scrim = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                           Inches(SW), Inches(SH))
            scrim.line.fill.background()
            spPr = scrim._element.spPr
            for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
                for el in spPr.findall(qn(tag)):
                    spPr.remove(el)
            g = spPr.makeelement(qn("a:gradFill"), {})
            lst = g.makeelement(qn("a:gsLst"), {})
            # 좌(진한 검정 70%) → 우(옅은 검정 22%)
            for pos, alpha in ((0, 72000), (62000, 42000), (100000, 20000)):
                gs = lst.makeelement(qn("a:gs"), {"pos": str(pos)})
                cc = gs.makeelement(qn("a:srgbClr"), {"val": "0E1420"})
                al = cc.makeelement(qn("a:alpha"), {"val": str(alpha)})
                cc.append(al); gs.append(cc); lst.append(gs)
            g.append(lst)
            g.append(g.makeelement(qn("a:lin"), {"ang": "0", "scaled": "1"}))
            spPr.append(g)
        except Exception:
            pass

    # 1) 좌측 세로 액센트 바(풀하이트)
    try:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                     Inches(0.46), Inches(SH))
        _grad(bar, ('%02X%02X%02X' % (_hex_to_rgb(primary) or (46, 91, 186))),
              ('%02X%02X%02X' % (_hex_to_rgb(secondary) or (91, 155, 213))), ang="5400000")
        bar.line.fill.background()
    except Exception:
        pass

    # 2) 우상단 기하 액센트 클러스터(대형 원 + 보조 원) — 색감/깊이 (이미지 위에선 생략)
    if not over_image:
        try:
            c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW - 2.7), Inches(-1.5),
                                        Inches(3.9), Inches(3.9))
            c1.fill.solid(); c1.fill.fore_color.rgb = _tint(primary, 0.82)
            c1.line.fill.background()
            c2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(SW - 1.5), Inches(0.7),
                                        Inches(1.7), Inches(1.7))
            c2.fill.solid(); c2.fill.fore_color.rgb = _tint(accent, 0.66)
            c2.line.fill.background()
        except Exception:
            pass

    # 텍스트 색 — 이미지 위에선 흰색 계열, 아니면 짙은 회색.
    _title_rgb = RGBColor(0xFF, 0xFF, 0xFF) if over_image else RGBColor(0x1A, 0x1A, 0x2A)
    _sub_rgb = RGBColor(0xE6, 0xEA, 0xF2) if over_image else RGBColor(0x55, 0x5C, 0x68)
    _eb_rgb = RGBColor(0xFF, 0xFF, 0xFF) if over_image else _rgb(primary)

    margin_x = 1.15
    # 3) 아이브로우 필 — 라이트 배경 pill + 좌측 컬러 룰 + 라벨 (레퍼런스 NOTICE 칩 스타일)
    eb_y = SH * 0.30
    eb_text = (eyebrow or "").strip()
    if eb_text:
        _ebl = eb_text.upper()[:60]
        _pill_w = min(SW - margin_x - 1.0, max(2.0, 0.95 + len(_ebl) * 0.135))
        if not over_image:
            try:
                pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(margin_x), Inches(eb_y - 0.16),
                                              Inches(_pill_w), Inches(0.5))
                pill.adjustments[0] = 0.5
                pill.fill.solid(); pill.fill.fore_color.rgb = _tint(primary, 0.90)
                pill.line.fill.background()
                _shadow(pill, alpha="14000", blur="60000")
            except Exception:
                pass
        try:
            rule = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(margin_x + 0.22), Inches(eb_y - 0.02),
                                          Inches(0.10), Inches(0.22))
            rule.adjustments[0] = 0.5
            rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(primary)
            rule.line.fill.background()
        except Exception:
            pass
        _txt(margin_x + 0.46, eb_y - 0.16, _pill_w - 0.5, 0.5,
             _ebl, 13, True, _eb_rgb,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, spacing=160)
    else:
        try:
            rule = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(margin_x), Inches(eb_y + 0.06),
                                          Inches(0.62), Inches(0.10))
            rule.adjustments[0] = 0.5
            rule.fill.solid(); rule.fill.fore_color.rgb = _rgb(primary)
            rule.line.fill.background()
        except Exception:
            pass

    # 4) 대제목(큰 글씨) — 길이에 따라 폰트 자동 축소 + 후반 키워드 컬러 하이라이트(doc2 스타일)
    title_y = eb_y + 0.55
    tlen = _text_units(title)
    title_pt = 46 if tlen <= 16 else (40 if tlen <= 26 else (34 if tlen <= 38 else 28))
    # 제목 박스의 실제 점유 높이 추정(폰트 pt·줄수 기반) — 고정 2.0 대신 사용(결함 A).
    # CJK 1글자 폭 ≈ 폰트 크기(pt→in), 행간 1.2. 최대 2.0 캡(기존 높이 초과 금지),
    # 짧은 제목 보존용 하한 0.9(부제 기존 위치를 max()로 유지하기 위함).
    _title_w_in = SW - margin_x - 1.0
    _char_w_in = max(0.01, title_pt / 72.0)
    _units_per_line = max(1.0, _title_w_in / _char_w_in)
    _title_lines = max(1, int(math.ceil(_text_units(title[:90]) / _units_per_line)))
    est_title_h = _title_lines * (title_pt * 1.2 / 72.0) + 0.10
    est_title_h = max(0.9, min(2.0, est_title_h))
    try:
        _ttb = slide.shapes.add_textbox(Inches(margin_x), Inches(title_y),
                                        Inches(SW - margin_x - 1.0), Inches(est_title_h))
        _ttf = _ttb.text_frame
        _ttf.word_wrap = True
        try:
            _ttf.vertical_anchor = MSO_ANCHOR.TOP
            _ttf.margin_left = Emu(0); _ttf.margin_right = Emu(0)
            _ttf.margin_top = Emu(0); _ttf.margin_bottom = Emu(0)
        except Exception:
            pass
        _tp = _ttf.paragraphs[0]
        _tp.alignment = PP_ALIGN.LEFT
        _hl_rgb = RGBColor(0x9D, 0xC3, 0xFF) if over_image else _rgb(secondary)
        _words = _strip_md(str(title[:90])).split()
        _runs = []
        if len(_words) >= 2:
            _split_at = len(_words) // 2  # 후반부 키워드 하이라이트
            _r1 = _tp.add_run(); _r1.text = " ".join(_words[:_split_at]) + " "
            _r1.font.color.rgb = _title_rgb; _runs.append(_r1)
            _r2 = _tp.add_run(); _r2.text = " ".join(_words[_split_at:])
            _r2.font.color.rgb = _hl_rgb; _runs.append(_r2)
        else:
            _r1 = _tp.add_run(); _r1.text = _strip_md(str(title[:90]))
            _r1.font.color.rgb = _title_rgb; _runs.append(_r1)
        for _r in _runs:
            _r.font.size = Pt(title_pt); _r.font.bold = True
            try:
                _r.font.name = "Malgun Gothic"
            except Exception:
                pass
    except Exception:
        _txt(margin_x, title_y, SW - margin_x - 1.0, est_title_h, title[:90],
             title_pt, True, _title_rgb,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # 5) 부제 — 제목 박스의 실제 점유 높이(est_title_h)를 반영해 비겹침으로 수직 스택.
    #    고정 오프셋(title_y + 1.05/0.9)을 desired 로 두되, vertical_stack 이
    #    title 박스 bottom + gap 이상으로 보장 → 결함 A(제목↔부제 세로 겹침) 제거.
    #    겹침이 없던 짧은 제목은 max(기존 sub_y, title_bottom+gap)로 기존 위치 보존.
    sub = (subtitle or "").strip()
    _orig_sub_y = title_y + (1.05 if title_pt >= 40 else 0.9)
    _title_rect = (margin_x, title_y, SW - margin_x - 1.0, est_title_h)
    _sub_rect_desired = (margin_x + 0.02, _orig_sub_y, SW - margin_x - 1.6, 1.0)
    _stacked = vertical_stack([_title_rect, _sub_rect_desired], gap=0.12)
    sub_y = _stacked[1][1]
    if sub:
        _txt(margin_x + 0.02, sub_y, SW - margin_x - 1.6, 1.0, sub[:140],
             19, False, _sub_rgb, align=PP_ALIGN.LEFT)

    # 6) 짧은 액센트 라인(부제 아래)
    line_y = sub_y + (0.72 if sub else 0.3)
    try:
        ln = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(margin_x), Inches(line_y),
                                    Inches(3.4), Inches(0.055))
        ln.adjustments[0] = 0.5
        ln.fill.solid(); ln.fill.fore_color.rgb = _rgb(accent)
        ln.line.fill.background()
    except Exception:
        pass

    drew = True

    # 7) 하단: KPI 카드 행(있으면) 또는 날짜 pill
    kpis = list(kpis or [])[:4]
    if kpis:
        n = len(kpis)
        gap = 0.3
        avail = SW - margin_x - 0.8
        card_w = (avail - (n - 1) * gap) / n
        card_h = 1.25
        cy = SH - card_h - 0.55
        for i, kp in enumerate(kpis):
            try:
                value, label, delta = (list(kp) + ["", "", ""])[:3]
            except Exception:
                value, label, delta = str(kp), "", ""
            cx = margin_x + i * (card_w + gap)
            sh_hex = pcols[i % len(pcols)] if pcols else primary
            try:
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              Inches(cx), Inches(cy),
                                              Inches(card_w), Inches(card_h))
                card.adjustments[0] = 0.10
                _grad(card, '%02X%02X%02X' % tuple(int(v + (255 - v) * 0.86)
                      for v in (_hex_to_rgb(sh_hex) or (46, 91, 186))), "FFFFFF")
                card.line.color.rgb = _rgb(sh_hex)
                card.line.width = Pt(1.25)
                _shadow(card)
                # 상단 색 띠
                cap = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                             Inches(cx), Inches(cy),
                                             Inches(card_w), Inches(0.09))
                cap.adjustments[0] = 0.5
                cap.fill.solid(); cap.fill.fore_color.rgb = _rgb(sh_hex)
                cap.line.fill.background()
            except Exception:
                pass
            _txt(cx + 0.12, cy + 0.16, card_w - 0.24, 0.55, str(value)[:14],
                 26, True, _rgb(sh_hex), align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
            if label:
                _txt(cx + 0.14, cy + 0.74, card_w - 0.28, 0.42, str(label)[:28],
                     11, False, RGBColor(0x55, 0x5C, 0x68), align=PP_ALIGN.LEFT)
    elif date_str:
        try:
            pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(margin_x), Inches(SH - 1.0),
                                          Inches(2.3), Inches(0.5))
            pill.adjustments[0] = 0.5
            pill.fill.solid(); pill.fill.fore_color.rgb = _tint(primary, 0.84)
            pill.line.color.rgb = _rgb(primary); pill.line.width = Pt(1.0)
            tf = pill.text_frame
            tf.word_wrap = True
            from pptx.enum.text import MSO_ANCHOR as _MA
            tf.vertical_anchor = _MA.MIDDLE
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(date_str)[:24]
            r.font.size = Pt(13); r.font.bold = True
            r.font.color.rgb = _rgb(primary)
            try:
                r.font.name = "Malgun Gothic"
            except Exception:
                pass
        except Exception:
            pass

    # 8) 우하단 코너 액센트 바
    try:
        cb = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    Inches(SW - 2.0), Inches(SH - 0.18),
                                    Inches(2.0), Inches(0.18))
        cb.fill.solid(); cb.fill.fore_color.rgb = _rgb(accent)
        cb.line.fill.background()
    except Exception:
        pass

    return drew
