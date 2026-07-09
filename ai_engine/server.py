"""FastAPI server — AI Editor backend."""
import os
import json
import uuid
import asyncio
import subprocess
import re
from collections import deque
from datetime import datetime

from fastapi import FastAPI, Request
from fastapi.responses import StreamingResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware

__version__ = "0.4.0"

app = FastAPI(title="AI Editor Engine", version=__version__)
app.add_middleware(CORSMiddleware, allow_origins=["*"], allow_methods=["*"], allow_headers=["*"])

# Startup banner — 서버가 새 코드로 실행 중인지 사용자가 즉시 확인 가능
print(f"[AI Editor Engine] v{__version__} loaded — deterministic merger + forced fallback active")

# ===== Agent Tool Definitions =====
AGENT_TOOLS = {
    "tools": [
        {
            "toolSpec": {
                "name": "read_file",
                "description": "파일 내용을 읽습니다. 프로젝트 내 모든 파일을 읽을 수 있습니다.",
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string", "description": "읽을 파일의 절대 경로 또는 프로젝트 상대 경로"}
                        },
                        "required": ["path"]
                    }
                }
            }
        },
        {
            "toolSpec": {
                "name": "write_file",
                "description": "파일에 내용을 씁니다. 새 파일 생성 또는 기존 파일 덮어쓰기.",
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string", "description": "쓸 파일의 절대 경로"},
                            "content": {"type": "string", "description": "파일에 쓸 내용"}
                        },
                        "required": ["path", "content"]
                    }
                }
            }
        },
        {
            "toolSpec": {
                "name": "list_directory",
                "description": "디렉토리의 파일/폴더 목록을 반환합니다.",
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "path": {"type": "string", "description": "탐색할 디렉토리 경로"}
                        },
                        "required": ["path"]
                    }
                }
            }
        },
        {
            "toolSpec": {
                "name": "run_command",
                "description": "터미널 명령어를 실행하고 결과를 반환합니다. git, npm, pip 등 모든 CLI 도구 사용 가능.",
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "command": {"type": "string", "description": "실행할 셸 명령어"},
                            "cwd": {"type": "string", "description": "작업 디렉토리 (선택)"}
                        },
                        "required": ["command"]
                    }
                }
            }
        },
        {
            "toolSpec": {
                "name": "search_files",
                "description": "프로젝트 내 파일에서 텍스트를 검색합니다 (grep).",
                "inputSchema": {
                    "json": {
                        "type": "object",
                        "properties": {
                            "query": {"type": "string", "description": "검색할 텍스트 또는 정규식"},
                            "path": {"type": "string", "description": "검색할 디렉토리 경로"},
                            "file_pattern": {"type": "string", "description": "파일 패턴 (예: *.py, *.js)"}
                        },
                        "required": ["query", "path"]
                    }
                }
            }
        }
                ,
            {
                "toolSpec": {
                    "name": "generate_image",
                    "description": "Generate an image from a text prompt using Bedrock image models. Saves PNG to .generated/. Supports models: Stability SD3.5, Stable Image Core, Amazon Titan Image v2.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "prompt": {"type": "string", "description": "Image description (max 2000 chars)"},
                                "size": {"type": "string", "description": "Image size like 1024x1024, 1024x768, 768x1024", "default": "1024x1024"},
                                "style": {"type": "string", "description": "Optional Stability style preset (photographic, cinematic, anime, etc.)"}
                            },
                            "required": ["prompt"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_pdf",
                    "description": "Generate a multi-page PDF document from structured sections. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Document title (cover page)"},
                                "sections": {
                                    "type": "array",
                                    "description": "List of sections, each with heading and body text",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "heading": {"type": "string"},
                                            "body": {"type": "string"}
                                        }
                                    }
                                }
                            },
                            "required": ["title", "sections"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_pptx",
                    "description": "Generate a PowerPoint presentation from structured slides. Each slide can have a title, bullets, and an optional imagePrompt to auto-generate an image. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Cover slide title"},
                                "slides": {
                                    "type": "array",
                                    "description": "List of slides",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "title": {"type": "string"},
                                            "bullets": {"type": "array", "items": {"type": "string"}},
                                            "imagePrompt": {"type": "string", "description": "Optional: auto-generate an image for the slide"},
                                            "layout": {"type": "string", "description": "title | content | two-column", "default": "content"}
                                        }
                                    }
                                }
                            },
                            "required": ["title", "slides"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "edit_image",
                    "description": "기존 이미지를 편집합니다. inpaint(마스크 영역 교체) 또는 outpaint(캔버스 확장) 모드를 지원합니다. Titan Image v2를 1순위, Nova Canvas를 폴백으로 사용합니다. inpaint 모드는 mask_path가 필요하고, outpaint 모드는 direction과 extend_pixels가 필요합니다.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "mode": {
                                    "type": "string",
                                    "enum": ["inpaint", "outpaint"],
                                    "description": "편집 모드: inpaint(마스크 영역 교체) 또는 outpaint(캔버스 확장)"
                                },
                                "image_path": {
                                    "type": "string",
                                    "description": "원본 이미지 경로 (PNG/JPEG, 최대 5MB)"
                                },
                                "prompt": {
                                    "type": "string",
                                    "minLength": 1,
                                    "maxLength": 1000,
                                    "description": "편집 프롬프트 (1~1000자)"
                                },
                                "mask_path": {
                                    "type": "string",
                                    "description": "마스크 이미지 경로 (inpaint 모드 필수, 흰색=편집 영역, 검정=보존 영역, 원본과 동일 해상도)"
                                },
                                "direction": {
                                    "type": "array",
                                    "items": {
                                        "type": "string",
                                        "enum": ["up", "down", "left", "right"]
                                    },
                                    "minItems": 1,
                                    "maxItems": 4,
                                    "description": "확장 방향 (outpaint 모드 필수, 1~4개)"
                                },
                                "extend_pixels": {
                                    "type": "integer",
                                    "minimum": 1,
                                    "maximum": 1024,
                                    "description": "확장 크기 픽셀 (outpaint 모드 필수, 1~1024)"
                                }
                            },
                            "required": ["mode", "image_path", "prompt"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_xlsx",
                    "description": "Generate an Excel workbook (XLSX) using openpyxl. Each sheet has headers (first row, bold) and rows of data. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "File title (used in slug + first-sheet header). e.g., 'sales-report'"},
                                "sheets": {
                                    "type": "array",
                                    "description": "List of sheets",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "name": {"type": "string", "description": "Sheet name (max 31 chars)"},
                                            "headers": {"type": "array", "items": {"type": "string"}, "description": "Header row (bold, accent fill)"},
                                            "rows": {
                                                "type": "array",
                                                "description": "Data rows — array of arrays. Cells can be string/number/bool.",
                                                "items": {"type": "array"}
                                            }
                                        }
                                    }
                                }
                            },
                            "required": ["title", "sheets"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_docx",
                    "description": "Generate a Word document (DOCX) using python-docx. Sections support headings (h1/h2/h3) and body paragraphs. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Document title (cover heading)"},
                                "sections": {
                                    "type": "array",
                                    "description": "List of sections",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "heading": {"type": "string"},
                                            "level": {"type": "integer", "description": "Heading level 1-3 (default 2)"},
                                            "body": {"type": "string", "description": "Section body. Newlines split into paragraphs."},
                                            "bullets": {"type": "array", "items": {"type": "string"}, "description": "Optional bullet list"}
                                        }
                                    }
                                }
                            },
                            "required": ["title", "sections"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_native_diagram",
                    "description": "Generate a native (editable-style) diagram PNG using matplotlib (no Bedrock call). diagram_type: tree (indented folder tree), flow (left-to-right arrows), architecture, stack, block. Saves PNG to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "diagram_type": {"type": "string", "description": "Diagram kind: tree, flow, architecture, stack, block"},
                                "title": {"type": "string", "description": "Diagram title"},
                                "content": {"type": "string", "description": "Text describing the diagram (lines for tree; 'A -> B -> C' for flow)"}
                            },
                            "required": ["diagram_type", "title", "content"]
                        }
                    }
                }
            }]
}


def _agent_tools_to_openai() -> list:
    """AGENT_TOOLS(Bedrock toolSpec) → OpenAI Responses 함수 도구 형식으로 변환.

    Bedrock:  {"toolSpec":{"name","description","inputSchema":{"json": <schema>}}}
    OpenAI:   {"type":"function","name","description","parameters": <schema>}
    게이트웨이 OpenAI Responses 라우트가 tools를 받으면 function_call을 반환한다.
    """
    out = []
    for t in AGENT_TOOLS.get("tools", []):
        spec = t.get("toolSpec") or {}
        name = spec.get("name")
        if not name:
            continue
        schema = ((spec.get("inputSchema") or {}).get("json")) or {"type": "object", "properties": {}}
        out.append({
            "type": "function",
            "name": name,
            "description": spec.get("description", ""),
            "parameters": schema,
        })
    return out


# ===== Remote Bridge Tool Routing =====
# When AE_BRIDGE_URL is set (by Electron), AI agent tools route through
# the bridge HTTP server which forwards to the remote SSH session.
import httpx as _httpx

_BRIDGE_URL = os.environ.get("AE_BRIDGE_URL", "")
_BRIDGE_TOKEN = os.environ.get("AE_BRIDGE_TOKEN", "")

# Dev mode fallback: read discovery file written by Electron main.js
if not _BRIDGE_URL:
    import tempfile as _tempfile
    try:
        _disc_path = os.path.join(_tempfile.gettempdir(), "ae-bridge.json")
        with open(_disc_path, "r") as _f:
            _disc = json.load(_f)
        _BRIDGE_URL = (_disc.get("url") or "").strip()
        _BRIDGE_TOKEN = (_disc.get("token") or "").strip()
    except (FileNotFoundError, json.JSONDecodeError, OSError):
        pass
_BRIDGE_TOOLS = {"read_file", "write_file", "list_directory", "search_files", "run_command"}

def _refresh_bridge_discovery():
    """Re-read bridge URL/token from env or discovery file (lazy).

    bridge_client.py와 *동일한* discovery 로직을 사용해 단일 진실 소스 유지.
    """
    global _BRIDGE_URL, _BRIDGE_TOKEN
    try:
        from ai_engine import bridge_client as _bc
    except ImportError:
        try:
            import bridge_client as _bc
        except ImportError:
            _bc = None
    if _bc is not None:
        _bc._load_discovery(force=True)
        if _bc._state.get("url") and _bc._state.get("token"):
            _BRIDGE_URL = _bc._state["url"]
            _BRIDGE_TOKEN = _bc._state["token"]
            return
    # bridge_client import 실패 시 직접 디스커버리
    url = os.environ.get("AE_BRIDGE_URL", "").strip()
    token = os.environ.get("AE_BRIDGE_TOKEN", "").strip()
    if url and token:
        _BRIDGE_URL, _BRIDGE_TOKEN = url, token
        return
    import tempfile as _tf
    try:
        with open(os.path.join(_tf.gettempdir(), "ae-bridge.json"), "r") as f:
            d = json.load(f)
        url = (d.get("url") or "").strip()
        token = (d.get("token") or "").strip()
        if url and token:
            _BRIDGE_URL, _BRIDGE_TOKEN = url, token
    except (FileNotFoundError, json.JSONDecodeError, OSError):
        pass


def _call_bridge(endpoint: str, payload: dict, timeout: float = 30.0):
    """Call the Electron bridge server. Returns dict or None on failure.

    매 호출마다 discovery를 한 번씩 refresh — Electron이 ai_engine보다 먼저
    뜨든 나중에 뜨든 자동으로 follow-up 가능. 캐시 유효시간 5초 내에는 no-op.
    """
    global _BRIDGE_URL, _BRIDGE_TOKEN
    # 매번 refresh (bridge_client._load_discovery에 5초 캐시 있음)
    _refresh_bridge_discovery()
    if not _BRIDGE_URL or not _BRIDGE_TOKEN:
        return None
    try:
        r = _httpx.post(
            f"{_BRIDGE_URL}/bridge/{endpoint}",
            headers={"X-AE-Bridge-Token": _BRIDGE_TOKEN},
            json=payload,
            timeout=timeout,
        )
        if r.status_code == 200:
            return r.json()
        if r.status_code in (401, 404, 502, 503):
            _BRIDGE_URL, _BRIDGE_TOKEN = "", ""
    except Exception:
        _BRIDGE_URL, _BRIDGE_TOKEN = "", ""
    return None

def _bridge_is_remote() -> bool:
    """Check if a remote SSH session is currently active."""
    res = _call_bridge("status", {})
    return bool(res and res.get("remote"))


def _format_bridge_result(tool_name, br):
    """Format bridge JSON response into a string for the agent."""
    if not isinstance(br, dict):
        return str(br)
    if not br.get('ok', True):
        return f'[Bridge Error] {br.get("error", "unknown")}'
    if tool_name == 'read_file':
        c = br.get('content', '')
        _max = int(os.environ.get('AE_READ_FILE_MAX', '120000'))
        if len(c) > _max:
            c = c[:_max] + '\n... (truncated)'
        return c
    elif tool_name == 'write_file':
        return 'File saved [remote]'
    elif tool_name == 'list_directory':
        entries = br.get('entries', [])
        lines = []
        for e in sorted(entries, key=lambda x: (not x.get('isDirectory'), x.get('name', ''))):
            nm = e.get('name', '')
            if nm.startswith('.') and nm not in ('.env', '.gitignore'):
                continue
            kind = 'DIR' if e.get('isDirectory') else 'FILE'
            lines.append(f'  {kind}  {nm}')
        return f'({len(lines)} items) [remote]\n' + '\n'.join(lines[:200])
    elif tool_name == 'run_command':
        output = (br.get('stdout', '') + br.get('stderr', ''))
        _max = int(os.environ.get('AE_RUN_CMD_MAX', '40000'))
        if len(output) > _max:
            output = output[:_max] + '\n... (truncated)'
        return output or '(no output)'
    elif tool_name == 'search_files':
        return br.get('output', 'no results')
    return str(br)



# ===== Media Generation Tools =====


def _resolve_relative_for_verify(rel_path: str, project_path: str = "") -> str:
    """`.generated/foo.pdf` 같은 상대경로를 도구가 실제 저장한 절대 경로로 매핑.

    `_tool_generate_pdf/pptx/...` 는 `_resolve_local_root(project_path)`로
    저장 위치를 결정하는데, 그 함수는 상황에 따라:
        1. project_path 자체 (작업 중인 폴더)
        2. AE_GENERATED_ROOT (Electron이 주입한 userData 경로)
        3. ~/.agentic-editor/
        4. tempdir
    중 하나를 반환한다. orchestrator의 검증 단계는 같은 후보 우선순위로
    파일을 찾아야 verified_files에 정확히 들어간다 (이전 버그: orchestrator가
    project_path로만 join해서 도구가 fallback 경로에 저장한 파일을 검증 단계
    에서 0건으로 판단 → 5명 에이전트 모두 실패 보고).

    절대경로면 그대로 반환. 후보를 모두 시도해 첫 존재 파일 절대경로 반환.
    아무 데서도 못 찾으면 첫 후보(project_path) 절대경로 반환 (caller가 isfile로
    재확인 후 실패 처리).
    """
    if os.path.isabs(rel_path):
        return rel_path
    candidates = []
    # 1) 작업 폴더
    if project_path and os.path.isdir(project_path):
        candidates.append(project_path)
    # 2) Electron 주입 userData
    env_root = os.environ.get("AE_GENERATED_ROOT", "").strip()
    if env_root:
        candidates.append(env_root)
    # 3) ~/.agentic-editor
    candidates.append(os.path.expanduser("~/.agentic-editor"))
    # 4) tempdir
    import tempfile as _tf
    candidates.append(_tf.gettempdir())
    # 5) cwd (최후 fallback — 일반적으로 위 후보 중 하나에 매칭됨)
    candidates.append(os.getcwd())

    for root in candidates:
        cand = os.path.join(root, rel_path)
        if os.path.isfile(cand):
            return cand
    # 못 찾음 — 첫 후보로 반환 (caller가 isfile 체크 후 실패 처리)
    return os.path.join(candidates[0], rel_path)


def _resolve_local_root(project_path: str = "") -> str:
    """파일 저장에 사용할 사용자별 쓰기 가능한 로컬 루트 경로 반환.

    사용자 정책 (30명 멀티유저 SSH 환경):
    - 원격 SSH 작업 폴더(/fsx/home/<user>/...)에는 절대 쓰지 않음 → 데이터 격리.
    - 항상 클라이언트 워크스테이션의 *로컬* 사용자 홈/userData에만 저장.
    - 결과물은 OS user별로 ~/.agentic-editor/.generated/ 에서만 접근 가능.

    우선순위:
    1. project_path가 *로컬*이고 쓰기 가능 → 그대로 사용 (e.g. /Users/jcg/myproject)
       - /fsx/, /home/ (Linux 원격), /opt/, FSx 마운트 등 원격 패턴은 1순위에서 제외.
    2. AE_GENERATED_ROOT 환경변수 (Electron이 userData 경로 주입)
    3. ~/.agentic-editor/ (사용자 홈, 항상 쓰기 가능, OS user별 격리)
    4. tempfile.gettempdir() (최후 fallback)

    30명 배포 시 핵심:
    - 사용자 A의 ~/.agentic-editor/ 와 사용자 B의 ~/.agentic-editor/는 OS 레벨에서 격리
    - 앱 설치 디렉토리(/Applications/...)는 읽기 전용이라 절대 사용 안 함
    - SSH 원격 작업 폴더는 공유 가능하므로 절대 미디어 파일 저장 위치로 사용 안 함
    """
    # 0) 원격 path 감지 — 사용자 정책상 원격엔 절대 쓰지 않음
    def _looks_remote(p: str) -> bool:
        if not p:
            return False
        # FSx 마운트, 일반적 SSH 원격 home/opt, [SSH:...] 마커, .ssh-mount 등
        if p.startswith("/fsx/") or p.startswith("/fsx") and len(p) > 4 and p[4] == "/":
            return True
        if p.startswith("/home/") or p.startswith("/opt/"):
            # macOS 사용자 홈은 /Users/이지 /home/이 아님 → /home/는 원격 신호.
            # Linux native에서 실행될 경우엔 AE_LOCAL_HOME=1 환경변수로 우회 가능.
            if os.environ.get("AE_LOCAL_HOME") == "1":
                return False
            return True
        if "[SSH:" in p:
            return True
        return False

    # 1) project_path 우선 — 단, 원격 패턴이면 skip
    if project_path and not _looks_remote(project_path) and os.path.isdir(project_path):
        try:
            # 쓰기 가능 확인
            test_file = os.path.join(project_path, ".ae_write_test")
            with open(test_file, "w") as f:
                f.write("")
            os.remove(test_file)
            return project_path
        except (OSError, PermissionError):
            pass

    # 2) Electron이 userData 경로를 환경변수로 주입한 경우 — 사용자별 OS 격리됨
    env_root = os.environ.get("AE_GENERATED_ROOT", "").strip()
    if env_root:
        try:
            os.makedirs(env_root, exist_ok=True)
            return env_root
        except (OSError, PermissionError):
            pass

    # 3) ~/.agentic-editor/ — 항상 사용자별로 OS 레벨 격리됨, 쓰기 가능
    home_root = os.path.expanduser("~/.agentic-editor")
    try:
        os.makedirs(home_root, exist_ok=True)
        return home_root
    except (OSError, PermissionError):
        pass

    # 4) tempdir fallback
    import tempfile as _tf
    return _tf.gettempdir()


# ===== Stability AI Image Services — Bedrock 등록 13+ 모델 =====
# 게이트웨이가 활성화한 경우 사용. 호출 시점 _resolve_callable_model_id가
# us./global. prefix 자동 부착.
STABILITY_GENERATIVE_IDS = [
    "stability.stable-image-ultra-v1:1",     # 사진 최고급
    "stability.sd3-5-large-v1:0",            # 범용 고품질
    "stability.stable-image-core-v1:1",      # 빠른 생성
]
STABILITY_INPAINT_IDS = [
    "stability.stable-image-inpaint-v1:0",
    "stability.stable-image-search-and-replace-v1:0",
]
STABILITY_OUTPAINT_IDS = [
    "stability.stable-image-outpaint-v1:0",
]
STABILITY_BACKGROUND_IDS = [
    "stability.stable-image-remove-background-v1:0",
]
STABILITY_ERASE_IDS = [
    "stability.stable-image-erase-object-v1:0",
]
STABILITY_RECOLOR_IDS = [
    "stability.stable-image-search-and-recolor-v1:0",
]
STABILITY_UPSCALE_IDS = [
    "stability.stable-image-creative-upscale-v1:0",
    "stability.stable-image-conservative-upscale-v1:0",
    "stability.stable-image-fast-upscale-v1:0",
]
STABILITY_CONTROL_IDS = [
    "stability.stable-image-control-sketch-v1:0",
    "stability.stable-image-control-structure-v1:0",
]
STABILITY_STYLE_IDS = [
    "stability.stable-image-style-transfer-v1:0",
    "stability.stable-image-style-guide-v1:0",
]

# Image generation model preference chain — ordered by general quality/recency.
# Stability 13개 + Amazon 2개를 모두 후보로 — generative 작업 fallback chain.
IMAGE_MODELS = (
    STABILITY_GENERATIVE_IDS
    + ["amazon.nova-canvas-v1:0", "amazon.titan-image-generator-v2:0"]
)

# 이미지 편집 — 작업별로 특화 모델이 명확히 분기됨.
IMAGE_EDIT_MODELS = (
    STABILITY_INPAINT_IDS
    + ["amazon.titan-image-generator-v2:0", "amazon.nova-canvas-v1:0"]
)
IMAGE_OUTPAINT_MODELS = (
    STABILITY_OUTPAINT_IDS
    + ["amazon.titan-image-generator-v2:0", "amazon.nova-canvas-v1:0"]
)
IMAGE_ERASE_MODELS = (
    STABILITY_ERASE_IDS + STABILITY_INPAINT_IDS
    + ["amazon.titan-image-generator-v2:0"]
)
IMAGE_BACKGROUND_REMOVE_MODELS = list(STABILITY_BACKGROUND_IDS)
IMAGE_RECOLOR_MODELS = STABILITY_RECOLOR_IDS + STABILITY_INPAINT_IDS
IMAGE_UPSCALE_MODELS = list(STABILITY_UPSCALE_IDS)
IMAGE_CONTROL_MODELS = list(STABILITY_CONTROL_IDS)
IMAGE_STYLE_MODELS = list(STABILITY_STYLE_IDS)


def _select_image_models(prompt: str, hint: str = "") -> list:
    """프롬프트 분석으로 이미지 모델 우선순위 결정.

    사용자가 "claude로 그려줘" 라고 해도 이 함수가 실제 이미지 생성 모델을
    프롬프트 키워드 기반으로 다시 선정합니다. 반환 리스트는 fallback 순서.

    Args:
        prompt: 이미지 생성 프롬프트
        hint: optional style/category hint (e.g., "photo", "diagram")

    Returns:
        list of model ids in preference order — first one is "best fit",
        나머지는 첫 번째 실패 시 fallback.
    """
    p = (prompt or "").lower()
    h = (hint or "").lower()
    text = p + " " + h

    # Photographic / 사진 / 리얼리즘 → Stability Ultra (최고 사진 품질)
    photo_kw = ("photo", "photograph", "realistic", "사진", "리얼리즘", "리얼한",
                "human", "portrait", "초상", "얼굴", "people", "사람",
                "cinematic", "영화", "studio", "현실적")
    # 다이어그램 / 차트 / UI / 스크린샷 → SD 3.5 Large (텍스트 렌더링 강함)
    diagram_kw = ("diagram", "chart", "flowchart", "architecture", "아키텍처",
                  "ui", "wireframe", "screenshot", "scheme", "schematic",
                  "다이어그램", "차트", "플로우차트", "구조도", "도식", "와이어프레임")
    # 일러스트 / 애니메이션 → Stability Ultra (스타일 다양)
    art_kw = ("illustration", "anime", "cartoon", "manga", "vector", "애니",
              "일러스트", "만화", "벡터", "그림체", "디지털 아트")
    # 로고 / 아이콘 / 단순 그래픽 → Stable Image Core (빠르고 저렴)
    logo_kw = ("logo", "icon", "emblem", "favicon", "아이콘", "로고",
               "심볼", "엠블럼", "minimal", "심플")
    # 기업 / 회의 / 비즈니스 자료 → Nova Canvas (AWS, 안전)
    biz_kw = ("business", "corporate", "presentation", "slide",
              "기업", "회사", "비즈니스", "회의", "보고서")

    def has_any(words):
        return any(w in text for w in words)

    # 우선순위 계산
    if has_any(diagram_kw):
        # 다이어그램 — SD 3.5 Large가 텍스트 렌더링이 가장 우수.
        # Titan은 게이트웨이 라우트가 일시적으로 500을 반환하는 상태이므로
        # chain에는 유지하되 우선순위 최하위로 둔다 (catalog에서는 절대 제거하지 않음).
        # Stable Image Core를 Nova Canvas 앞에 배치 — Core는 안정적으로 작동.
        return [
            "stability.sd3-5-large-v1:0",
            "stability.stable-image-ultra-v1:1",
            "stability.stable-image-core-v1:1",
            "amazon.nova-canvas-v1:0",
            "amazon.titan-image-generator-v2:0",
        ]
    if has_any(photo_kw):
        # 사진 — Ultra가 최고 사진 품질
        return [
            "stability.stable-image-ultra-v1:1",
            "stability.sd3-5-large-v1:0",
            "amazon.nova-canvas-v1:0",
            "amazon.titan-image-generator-v2:0",
        ]
    if has_any(logo_kw):
        # 로고 — Core가 빠르고 단순한 그래픽에 적합
        return [
            "stability.stable-image-core-v1:1",
            "stability.sd3-5-large-v1:0",
            "amazon.titan-image-generator-v2:0",
        ]
    if has_any(art_kw):
        # 일러스트 — Ultra의 스타일 다양성
        return [
            "stability.stable-image-ultra-v1:1",
            "stability.sd3-5-large-v1:0",
            "amazon.nova-canvas-v1:0",
        ]
    if has_any(biz_kw):
        # 비즈니스 자료 — Nova Canvas (안전, 보수적)
        return [
            "amazon.nova-canvas-v1:0",
            "stability.sd3-5-large-v1:0",
            "stability.stable-image-ultra-v1:1",
        ]
    # 기본 — 일반 fallback chain
    return list(IMAGE_MODELS)


# ===== Image Generation Circuit Breaker =====
# Bedrock 이미지 게이트웨이가 access denied(IAM/route)로 차단되면, 매 섹션마다
# 5개 모델 fallback × 60s timeout이 누적되어 전체 작업이 30분+로 폭주한다.
# 한번 실패가 감지되면 일정 시간(TTL) 동안 모든 호출을 즉시 단락(short-circuit)시켜
# PDF/PPTX/DOCX 생성을 텍스트/네이티브-다이어그램 경로로 빠르게 떨어뜨린다.
_IMAGE_GEN_CIRCUIT = {"disabled_at": 0, "ttl": 300}
# Access-denied/route-missing 패턴 — Bedrock 게이트웨이 거부 응답에 등장
_IMAGE_GEN_DENY_PATTERNS = (
    "access denied",
    "accessdenied",
    "execute-api:invoke",
    "principal identity",
    "http 403",
    "http 404",
    "not authorized",
    "forbidden",
)

# Ring buffer of recent image-gen attempts (last 10).
# Single source of truth for the diagnostic endpoint (`/api/debug/image-gen-status`)
# and for the short-circuit response's `recentAttempts` / `actionable` fields.
# Spec: media-output-quality (bugfix) — Property 2 / Req 2.3, 2.4, 2.5.
_IMAGE_GEN_ATTEMPTS: "deque[dict]" = deque(maxlen=10)


def _record_image_attempt(model: str, status: str, reason: str, duration_ms: int) -> None:
    """Append one image-gen attempt outcome to the ring buffer.

    Called from `_tool_generate_image` after the parallel `asyncio.gather`
    resolves. `status ∈ {"ok", "error", "exception"}`. `reason` is the
    gateway detail string for errors / repr(exc) for exceptions / empty
    string for successes.
    """
    import time as _t
    _IMAGE_GEN_ATTEMPTS.append({
        "ts": _t.time(),
        "model": model,
        "status": status,
        "reason": reason,
        "durationMs": duration_ms,
    })


def _image_gen_is_circuit_broken() -> bool:
    """Image gen이 access denied로 차단됐는지 확인. TTL 5분 후 자동 복구."""
    import time as _t
    if _IMAGE_GEN_CIRCUIT["disabled_at"] == 0:
        return False
    if _t.time() - _IMAGE_GEN_CIRCUIT["disabled_at"] > _IMAGE_GEN_CIRCUIT["ttl"]:
        # TTL expired → reset and try once more
        _IMAGE_GEN_CIRCUIT["disabled_at"] = 0
        return False
    return True


def _image_gen_trip_circuit(reason: str = ""):
    """Access denied 감지 → 회로 차단."""
    import time as _t
    _IMAGE_GEN_CIRCUIT["disabled_at"] = _t.time()
    print(f"[ImageGen Circuit] 차단됨 — {reason[:200]} (TTL {_IMAGE_GEN_CIRCUIT['ttl']}초)")


def _image_gen_error_is_access_denied(detail: str) -> bool:
    """에러 detail이 게이트웨이 access-denied 패턴이면 True."""
    if not detail:
        return False
    low = detail.lower()
    return any(pat in low for pat in _IMAGE_GEN_DENY_PATTERNS)


# ===== Korean Font Resolution =====
# reportlab과 python-docx는 기본 폰트(Helvetica/Calibri)에 한글 글리프가 없어
# PDF/DOCX에 한글이 ■■■ 박스로 깨진다. macOS/Linux/Windows별 시스템 폰트를
# 탐색하여 첫 발견 폰트를 모듈 단위로 캐시한다. matplotlib도 같은 폰트를 사용.
_KOREAN_FONT_CACHE = {"path": None, "name": None, "registered": False, "checked": False}

# macOS/Linux/Windows 후보 — 일반적으로 설치되어 있는 한글 폰트 우선순위.
# reportlab은 PostScript outline TTC를 지원 안 하므로 .ttf를 .ttc보다 앞에 둠.
_KOREAN_FONT_CANDIDATES = [
    # macOS — AppleGothic.ttf (Supplemental에 있고 reportlab 호환)
    ("/System/Library/Fonts/Supplemental/AppleGothic.ttf", "AppleGothic"),
    ("/Library/Fonts/AppleGothic.ttf", "AppleGothic"),
    # Linux — Nanum/Noto (.ttf)
    ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", "NanumGothic"),
    ("/usr/share/fonts/nanum/NanumGothic.ttf", "NanumGothic"),
    ("/usr/share/fonts/opentype/noto/NotoSansKR-Regular.otf", "NotoSansKR"),
    # Windows — Malgun (.ttf)
    ("C:/Windows/Fonts/malgun.ttf", "MalgunGothic"),
    ("C:/Windows/Fonts/NanumGothic.ttf", "NanumGothic"),
    # macOS — Noto OTF/CJK가 있으면 사용
    ("/Library/Fonts/NotoSansCJKkr-Regular.otf", "NotoSansCJKkr"),
    ("/Library/Fonts/NotoSansKR-Regular.otf", "NotoSansKR"),
    # Linux Noto CJK collections (.ttc)
    ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", "NotoSansCJK"),
    ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", "NotoSansCJK"),
    # macOS — AppleSDGothicNeo .ttc (마지막 — reportlab은 지원 안 하지만 matplotlib은 OK)
    ("/System/Library/Fonts/AppleSDGothicNeo.ttc", "AppleSDGothicNeo"),
]


def _resolve_korean_font_path() -> tuple:
    """시스템에서 한글 폰트 파일을 찾아 (path, name) 반환. 없으면 (None, None).

    파일 단위로 한 번만 탐색 후 캐시. AE_KOREAN_FONT 환경변수가 있으면 1순위.
    """
    if _KOREAN_FONT_CACHE["checked"]:
        return _KOREAN_FONT_CACHE["path"], _KOREAN_FONT_CACHE["name"]
    _KOREAN_FONT_CACHE["checked"] = True

    env_path = os.environ.get("AE_KOREAN_FONT", "").strip()
    if env_path and os.path.isfile(env_path):
        name = os.path.splitext(os.path.basename(env_path))[0].replace(" ", "")
        _KOREAN_FONT_CACHE["path"] = env_path
        _KOREAN_FONT_CACHE["name"] = name
        return env_path, name

    for path, name in _KOREAN_FONT_CANDIDATES:
        if os.path.isfile(path):
            _KOREAN_FONT_CACHE["path"] = path
            _KOREAN_FONT_CACHE["name"] = name
            print(f"[Korean Font] resolved: {name} ({path})")
            return path, name
    print("[Korean Font] 시스템에서 한글 폰트를 찾지 못함 — 한글 깨질 수 있음")
    return None, None


def _register_korean_font_for_reportlab() -> str:
    """reportlab에 한글 폰트 등록. 후보를 순회하며 첫 등록 성공한 폰트 사용.

    Apple SD Gothic Neo .ttc는 PostScript outline 기반이라 reportlab이 거부하므로
    실패 시 다음 후보(.ttf)를 시도한다. 모든 후보가 실패하면 Helvetica fallback.

    Returns:
        등록된 한글 폰트 이름 또는 "Helvetica" (전체 실패 시).
    """
    if _KOREAN_FONT_CACHE["registered"] and _KOREAN_FONT_CACHE["name"]:
        return _KOREAN_FONT_CACHE["name"]
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
    except ImportError:
        return "Helvetica"

    # AE_KOREAN_FONT 환경변수 우선
    candidates = []
    env_path = os.environ.get("AE_KOREAN_FONT", "").strip()
    if env_path and os.path.isfile(env_path):
        env_name = os.path.splitext(os.path.basename(env_path))[0].replace(" ", "")
        candidates.append((env_path, env_name))
    candidates.extend([(p, n) for p, n in _KOREAN_FONT_CANDIDATES if os.path.isfile(p)])

    for path, name in candidates:
        try:
            try:
                pdfmetrics.registerFont(TTFont(name, path))
            except TypeError:
                pdfmetrics.registerFont(TTFont(name, path, subfontIndex=0))
            try:
                from reportlab.pdfbase.pdfmetrics import registerFontFamily
                registerFontFamily(name, normal=name, bold=name, italic=name, boldItalic=name)
            except Exception:
                pass
            _KOREAN_FONT_CACHE["registered"] = True
            _KOREAN_FONT_CACHE["name"] = name
            _KOREAN_FONT_CACHE["path"] = path
            print(f"[Korean Font] reportlab 등록 성공: {name} ({path})")
            return name
        except Exception as e:
            print(f"[Korean Font] reportlab '{name}' 등록 실패 ({e}) — 다음 후보 시도")
            continue
    print("[Korean Font] 모든 후보 실패 — Helvetica fallback (한글 깨짐 가능)")
    return "Helvetica"


def _apply_korean_font_for_matplotlib():
    """matplotlib 전역 한글 폰트 적용. 한 번만 등록되도록 캐시 활용."""
    path, name = _resolve_korean_font_path()
    if not path:
        return
    try:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        from matplotlib import font_manager as _fm
        # font 파일을 직접 등록 (시스템에 fc-cache 안 된 환경 대응)
        try:
            _fm.fontManager.addfont(path)
        except Exception:
            pass
        try:
            prop = _fm.FontProperties(fname=path)
            plt.rcParams["font.family"] = prop.get_name()
        except Exception:
            plt.rcParams["font.family"] = name
        plt.rcParams["axes.unicode_minus"] = False
    except Exception as e:
        print(f"[Korean Font] matplotlib 적용 실패: {e}")


async def _try_vertex_image_single(prompt, size, project_path, aws_profile="", style=""):
    """Vertex AI(Nano Banana Pro / Imagen) 단일 이미지 폴백 — 최고 품질 경로.

    Bedrock 이미지 모델이 불가용(서킷 차단/403/전부 실패)일 때 호출된다.
    steering(gateway.md)의 이미지 예외 정책에 따라 *이미지 생성 한정*으로 Vertex 허용.
    Vertex 키가 해석되지 않거나(미설정) 실패하면 None을 반환해 기존 동작을 보존한다.
    반환: _tool_generate_image와 동일 형태의 payload dict 또는 None.
    """
    try:
        from ai_engine.vertex_image_module import get_vertex_image_client
        client = get_vertex_image_client(aws_profile=aws_profile or "")
    except Exception as _e:
        print(f"[ImageGen/Vertex] probe 실패: {str(_e)[:160]}")
        return None
    if not client or not getattr(client, "enabled", False):
        return None
    try:
        w, h = (int(x) for x in str(size).lower().split("x"))
    except Exception:
        w, h = 1024, 1024
    # Vertex는 정해진 종횡비 토큰만 허용한다 — "1024:768" 같은 픽셀값을 그대로
    # 넘기면 http-400. 픽셀 크기를 가장 가까운 지원 비율로 매핑한다.
    _ar = (float(w) / float(h)) if h else 1.0
    _AR_CHOICES = (("16:9", 16 / 9), ("4:3", 4 / 3), ("1:1", 1.0),
                   ("3:4", 3 / 4), ("9:16", 9 / 16))
    aspect = min(_AR_CHOICES, key=lambda c: abs(c[1] - _ar))[0]
    model_class = (os.environ.get("AE_VERTEX_MODEL_CLASS", "").strip()
                   or "image_generation_high_quality")
    _full_prompt = prompt
    if style:
        _full_prompt = f"{prompt}\n\nStyle: {style}"
    try:
        vres = await client.generate(
            prompt=_full_prompt, model_class=model_class, aspect_ratio=aspect,
            negative_prompt="watermark, fake logo, brand name, distorted UI, unreadable artifacts",
        )
    except Exception as _e:
        print(f"[ImageGen/Vertex] generate 예외: {str(_e)[:160]}")
        return None
    if not isinstance(vres, dict) or vres.get("error"):
        print(f"[ImageGen/Vertex] error: {str((vres or {}).get('error'))[:160]}")
        return None
    images = vres.get("images") or []
    if not images:
        return None
    import time as _t, hashlib as _hh, base64 as _b64
    local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(local_root, ".generated")
    try:
        os.makedirs(gen_dir, exist_ok=True)
    except OSError:
        return None
    ts = str(int(_t.time() * 1000))
    short = _hh.md5(_full_prompt.encode()).hexdigest()[:6]
    model_v = (vres.get("model") or "vertex").split("/")[-1].replace("-", "")[:24]
    fname = f"vertex-{ts}-{short}-{model_v}.png"
    abs_path = os.path.join(gen_dir, fname)
    try:
        with open(abs_path, "wb") as _f:
            _f.write(_b64.b64decode(images[0]))
        sz = os.path.getsize(abs_path)
        if sz <= 0:
            return None
    except Exception as _e:
        print(f"[ImageGen/Vertex] save 실패: {str(_e)[:160]}")
        return None
    aw, ah = w, h
    try:
        from PIL import Image as _PIL
        with _PIL.open(abs_path) as im:
            aw, ah = im.size
    except Exception:
        pass
    print(f"[ImageGen] Vertex 폴백 성공 — {fname} ({sz} bytes, {aw}x{ah})")
    return {
        "path": f".generated/{fname}", "absPath": abs_path,
        "model": vres.get("model") or "vertex-image",
        "size": f"{w}x{h}", "width": aw, "height": ah,
        "sizeBytes": sz, "qualityScore": 95, "via": "vertex",
    }


async def _tool_generate_image(tool_input: dict, project_path: str, aws_profile: str = '', bedrock_user: str = '') -> str:  # [patched-credentials]
    """Generate an image via Bedrock image models — *병렬 best-of-N + 품질 스코어*.

    이전 sequential fallback 방식은 첫 번째 성공한 모델 결과를 그대로 반환했고,
    한 모델이 실패하면 다음 모델로 넘어가는 데 60초씩 걸렸다 (5개 모델 × 60초).
    품질 보장 없이 가장 빨리 응답한 모델 = 가장 좋은 모델이 아님.

    이 함수는:
    1. 프롬프트 의도에 맞는 상위 N개 이미지 모델을 *병렬* 호출 (asyncio.gather)
    2. 성공한 결과들을 품질 스코어로 평가
       - 파일 크기 (너무 작으면 비어있을 가능성)
       - 해상도 정확도
       - Pillow entropy (디테일/복잡도)
       - 모델 우선순위 (선호 모델 가산점)
    3. 가장 높은 점수의 결과를 최종 반환, 나머지는 디스크에 보존 (사용자가 비교 가능)
    4. 회로 차단은 *모든* 호출이 access-denied일 때만 발동 (단일 실패엔 안 발동)
    5. AE_IMAGE_PARALLEL_N 환경변수로 병렬 호출 수 조정 (기본 5, 최대 5)

    Returns JSON string: {path, model, width, height, sizeBytes, qualityScore, candidates}
    candidates는 다른 모델이 만든 이미지들의 메타 (사용자가 패널에서 다른 결과 비교 가능).
    """
    import time as _t, hashlib, base64
    prompt = (tool_input.get("prompt") or "").strip()
    size = tool_input.get("size", "1024x1024")
    style = tool_input.get("style", "")

    if not prompt:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt is required"})
    if len(prompt) > 2000:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt exceeds 2000 chars"})

    # === 최고 품질 우선 — Vertex(Nano Banana Pro / Imagen 4) 1순위 ===
    # 채팅 모델이 무엇이든(레거시 GPT 5.4 포함) 이미지는 시스템이 Vertex로 생성한다.
    # Bedrock 이미지(텍스트 렌더 약함 + 현재 403/서킷)를 우회해 시간·품질을 모두 확보.
    # Vertex 키 미해석(미활성) 또는 실패 시 None → 아래 Bedrock 경로로 자연 폴백.
    # AE_PREFER_VERTEX_IMAGE=0 으로 비활성 가능. 재시도(_isRetry) 호출은 제외.
    if os.environ.get("AE_PREFER_VERTEX_IMAGE", "1") == "1" and not tool_input.get("_isRetry"):
        _vx_first = await _try_vertex_image_single(prompt, size, project_path, aws_profile, style)
        if _vx_first:
            print("[ImageGen] Vertex 1순위 생성 성공(고품질)")
            return json.dumps(_vx_first)

    if _image_gen_is_circuit_broken():
        # 회로 차단 — Bedrock 이미지 전부 거부 상태. 최고 품질 Vertex 폴백 우선 시도.
        _vx = await _try_vertex_image_single(prompt, size, project_path, aws_profile, style)
        if _vx:
            return json.dumps(_vx)
        import time as _t_local
        _disabled_at = _IMAGE_GEN_CIRCUIT.get("disabled_at", 0)
        _ttl = _IMAGE_GEN_CIRCUIT.get("ttl", 300)
        _remaining = max(0, int(_ttl - (_t_local.time() - _disabled_at)))
        msg = (
            f"이미지 모델 회로 차단 — TTL {_remaining}초 남음. "
            "원인: 이전 이미지 모델 호출 모두 access-denied. "
            "matplotlib fallback로만 동작. 게이트웨이 IAM 권한 점검 필요."
        )
        print(f"[ImageGen] {msg}")
        # Enrich the short-circuit payload with last-5 attempts and a bilingual
        # `actionable` message that names the offending model ids — so the chat
        # panel can show "Bedrock 게이트웨이가 image-gen 라우트를 거부했습니다 —
        # 권한 필요 모델: [...]" instead of a silent matplotlib fallback.
        # Spec: media-output-quality (bugfix) — Req 2.3, 2.5.
        recent_attempts = list(_IMAGE_GEN_ATTEMPTS)[-5:]
        denied_ids = sorted({
            a["model"]
            for a in recent_attempts
            if _image_gen_error_is_access_denied(a.get("reason", ""))
        })
        payload = {
            "error": "circuit-breaker",
            "detail": msg,
            "circuitTtlRemainingSec": _remaining,
            "recentAttempts": recent_attempts,
        }
        if denied_ids:
            ids_str = ", ".join(denied_ids)
            payload["actionable"] = (
                f"Bedrock 게이트웨이가 image-gen 라우트를 거부했습니다 — "
                f"권한 필요 모델: {ids_str} / "
                f"Bedrock gateway denied image-gen route — admin must grant "
                f"invoke permission for: {ids_str}"
            )
        return json.dumps(payload)

    try:
        w, h = (int(x) for x in size.lower().split("x"))
    except Exception:
        w, h = 1024, 1024

    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    ts = str(int(_t.time() * 1000))
    short_hash = hashlib.md5(prompt.encode()).hexdigest()[:4]

    aws_profile = aws_profile or os.environ.get("AWS_PROFILE", "bedrock-gw")
    bedrock_user = bedrock_user or os.environ.get("BEDROCK_USER", "")
    gw = _get_gw(aws_profile, bedrock_user)

    # 병렬 호출 — 프롬프트 분석 결과 상위 N개 모델 (기본 5개, env로 조정)
    # 5개로 한 이유: chain에 5개 모델(Stability 3 + Nova Canvas + Titan)이 있고,
    # Titan이 일시적으로 500을 반환해도 다른 4개가 병렬로 처리되어 막히지 않음.
    parallel_n = max(1, min(5, int(os.environ.get("AE_IMAGE_PARALLEL_N", "5"))))
    selected_models = _select_image_models(prompt, hint=style)[:parallel_n]
    print(f"[ImageGen] 병렬 호출 — {len(selected_models)}개 모델: {[m.split('.')[1] for m in selected_models]}")

    def _build_body(model_id: str) -> dict:
        if model_id.startswith("stability."):
            # Stability 는 픽셀비가 아니라 고정 enum(1:1, 16:9, 9:16 등)만 허용.
            # f"{w}:{h}"(예 '1024:768')는 ValidationException 을 유발하므로
            # 가로/세로/정사각을 유효 enum 으로 매핑한다.
            _ar = "1:1" if w == h else ("16:9" if w > h else "9:16")
            body = {
                "prompt": prompt,
                "mode": "text-to-image",
                "output_format": "png",
                "aspect_ratio": _ar,
            }
            if style:
                body["style_preset"] = style
            return body
        if model_id.startswith("amazon.titan"):
            return {
                "textToImageParams": {"text": prompt},
                "imageGenerationConfig": {
                    "numberOfImages": 1,
                    "width": w,
                    "height": h,
                    "quality": "premium",  # standard → premium 으로 품질 향상
                },
            }
        return {"prompt": prompt, "width": w, "height": h}

    def _save_and_score(idx: int, model_id: str, img_b64: str) -> dict:
        """디스크 저장 + 품질 스코어링."""
        try:
            img_bytes = base64.b64decode(img_b64)
        except Exception as _e:
            return {"ok": False, "model": model_id, "reason": f"base64 decode 실패: {_e}"}
        if not img_bytes or len(img_bytes) < 5_000:
            return {"ok": False, "model": model_id, "reason": f"PNG 너무 작음 ({len(img_bytes)} bytes)"}
        # 모델별 식별 가능한 파일명 (사용자가 패널에서 비교 가능)
        model_short = model_id.split(".")[1].split("-")[0]
        filename = f"image-{ts}-{short_hash}-{model_short}-{idx}.png"
        out_path = os.path.join(gen_dir, filename)
        with open(out_path, "wb") as f:
            f.write(img_bytes)
        # dimensions + entropy
        aw_, ah_, entropy = w, h, 0.0
        try:
            from PIL import Image as _PIL
            with _PIL.open(out_path) as im:
                aw_, ah_ = im.size
                # entropy: 이미지가 단조롭지 않은지 (0~8 범위, 사진은 보통 6~7)
                try:
                    entropy = float(im.convert("L").entropy())
                except Exception:
                    entropy = 0.0
        except Exception:
            pass
        # 품질 스코어 — 0~100
        # - 파일 크기 보너스 (너무 작으면 빈 이미지 의심)
        # - 해상도 정확도
        # - 엔트로피 (디테일)
        # - 모델 우선순위 (selected_models 순서가 의도 매칭도)
        size_score = min(40, len(img_bytes) / 10_000)  # 400KB → 40점 cap
        res_score = 20 if (aw_ == w and ah_ == h) else max(0, 20 - abs(aw_ - w) // 50)
        ent_score = min(25, max(0, (entropy - 3.0) * 10))  # entropy 5.5 → 25점
        try:
            priority_score = max(0, 15 - selected_models.index(model_id) * 3)
        except ValueError:
            priority_score = 0
        quality = round(size_score + res_score + ent_score + priority_score, 1)
        return {
            "ok": True,
            "path": f".generated/{filename}",
            "absPath": out_path,
            "model": model_id,
            "size": f"{w}x{h}",
            "width": aw_,
            "height": ah_,
            "sizeBytes": len(img_bytes),
            "entropy": round(entropy, 2),
            "qualityScore": quality,
        }

    async def _try_one(idx: int, model_id: str) -> dict:
        try:
            body = _build_body(model_id)
            callable_id = _resolve_callable_model_id(model_id, aws_profile, bedrock_user)
            result = await gw.invoke_model(callable_id, body, timeout=60)
            if "error" in result:
                return {"ok": False, "model": model_id, "reason": result["error"][:200]}
            images = result.get("images", [])
            if not images and isinstance(result, dict):
                images = result.get("artifacts", [])
                if images and isinstance(images[0], dict):
                    images = [a.get("base64", "") for a in images]
            if not images:
                return {"ok": False, "model": model_id, "reason": "no images returned"}
            img_b64 = images[0] if isinstance(images[0], str) else (
                images[0].get("base64", "") if isinstance(images[0], dict) else "")
            if not img_b64:
                return {"ok": False, "model": model_id, "reason": "empty image data"}
            return _save_and_score(idx, model_id, img_b64)
        except Exception as e:
            return {"ok": False, "model": model_id, "reason": f"exception: {str(e)[:200]}"}

    # 병렬 실행
    tasks = [_try_one(i, mid) for i, mid in enumerate(selected_models)]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    # Record every attempt outcome to the ring buffer so the diagnostic
    # endpoint and short-circuit response can surface real causes.
    # Spec: media-output-quality (bugfix) — Req 2.3, 2.4.
    for _i, _model_id in enumerate(selected_models):
        _r = results[_i] if _i < len(results) else None
        if isinstance(_r, BaseException):
            _record_image_attempt(_model_id, "exception", repr(_r)[:200], 0)
        elif isinstance(_r, dict) and _r.get("ok"):
            _record_image_attempt(_model_id, "ok", "", 0)
        elif isinstance(_r, dict):
            _record_image_attempt(_model_id, "error", (_r.get("reason") or "")[:200], 0)
        else:
            _record_image_attempt(_model_id, "error", f"unexpected result type: {type(_r).__name__}", 0)

    successes = []
    failures = []
    all_denied = True
    for r in results:
        if isinstance(r, Exception):
            failures.append({"model": "?", "reason": f"task exception: {r}"})
            continue
        if r.get("ok"):
            successes.append(r)
            all_denied = False
        else:
            failures.append(r)
            reason = (r.get("reason") or "").lower()
            if not _image_gen_error_is_access_denied(reason):
                all_denied = False

    # 모두 실패 + 모두 access-denied → 회로 차단
    if not successes:
        if all_denied and failures:
            _image_gen_trip_circuit(f"all {len(failures)} models access-denied")
        # Bedrock 전부 실패 → 최고 품질 Vertex 폴백 시도(가용 시).
        _vx = await _try_vertex_image_single(prompt, size, project_path, aws_profile, style)
        if _vx:
            return json.dumps(_vx)
        last_detail = failures[-1].get("reason", "all image models failed") if failures else "no models tried"
        hint = ""
        if any(t in last_detail for t in ("execute-api:Invoke", "principal identity", "HTTP 403", "HTTP 404")):
            hint = "현재 게이트웨이가 이미지 생성 라우트(/invoke-model)를 지원하지 않습니다. 관리자에게 활성화를 요청하세요."
        payload = {
            "error": "model-unavailable",
            "detail": last_detail[:200],
            "attempts": [{"model": f.get("model"), "reason": (f.get("reason") or "")[:120]} for f in failures],
        }
        if hint:
            payload["hint"] = hint
        return json.dumps(payload)

    # 성공한 결과들 중 가장 높은 품질 스코어 선택
    successes.sort(key=lambda r: r.get("qualityScore", 0), reverse=True)
    winner = successes[0]
    print(
        f"[ImageGen] best-of-{len(successes)} 선택 — winner={winner['model']} score={winner['qualityScore']}"
        f" (다른 후보: {[(s['model'].split('.')[1], s['qualityScore']) for s in successes[1:]]})"
    )

    # === 품질 임계값 미달 시 1회 자동 재생성 (개선된 프롬프트로) ===
    # 환경변수 AE_IMAGE_QUALITY_THRESHOLD (기본 60). 회귀 모드에서는 0으로 끔.
    # 재생성은 상위 1개 모델만, 더 디테일한 프롬프트로 시도 — 너무 비싸지지 않게.
    quality_threshold = float(os.environ.get("AE_IMAGE_QUALITY_THRESHOLD", "60"))
    is_retry = bool(tool_input.get("_isRetry"))
    if winner.get("qualityScore", 0) < quality_threshold and not is_retry:
        print(f"[ImageGen] 품질 임계 {quality_threshold} 미달 ({winner['qualityScore']}) — 강화 프롬프트로 1회 재시도")
        enhanced_prompt = (
            prompt
            + ", high detail, sharp focus, professional, vibrant colors, "
            + "8k uhd quality, well-composed, balanced lighting"
        )
        retry_input = dict(tool_input)
        retry_input["prompt"] = enhanced_prompt[:2000]
        retry_input["_isRetry"] = True  # 무한 재귀 방지
        retry_result = await _tool_generate_image(
            retry_input, project_path,
            aws_profile=aws_profile, bedrock_user=bedrock_user,
        )
        try:
            rp = json.loads(retry_result)
            if "error" not in rp and rp.get("qualityScore", 0) > winner.get("qualityScore", 0):
                print(f"[ImageGen] 재생성이 더 좋음 ({rp['qualityScore']} > {winner['qualityScore']}) — 재생성 결과 채택")
                # 재생성 후보들 + 원본 winner를 candidates에 합쳐서 반환
                merged_candidates = [
                    {
                        "path": winner["path"], "model": winner["model"],
                        "qualityScore": winner["qualityScore"], "sizeBytes": winner["sizeBytes"],
                    }
                ] + (rp.get("candidates") or [])
                rp["candidates"] = merged_candidates[:5]
                rp["retried"] = True
                return json.dumps(rp)
            else:
                print(f"[ImageGen] 재생성 결과가 더 안 좋거나 같음 — 원본 winner 유지")
        except (json.JSONDecodeError, TypeError):
            pass

    # winner를 메인으로 반환, 다른 후보는 candidates 메타로
    payload = {
        "path": winner["path"],
        "absPath": winner.get("absPath", ""),  # TASK 8 근본수정 — 카드 다운로드 정확성
        "model": winner["model"],
        "size": winner["size"],
        "width": winner["width"],
        "height": winner["height"],
        "sizeBytes": winner["sizeBytes"],
        "qualityScore": winner["qualityScore"],
        "entropy": winner.get("entropy", 0),
        "candidates": [
            {
                "path": s["path"], "model": s["model"],
                "qualityScore": s["qualityScore"], "sizeBytes": s["sizeBytes"],
            }
            for s in successes[1:]
        ],
        "attemptsTotal": len(selected_models),
        "successCount": len(successes),
    }
    return json.dumps(payload)


# ===== Native Diagram Generator (matplotlib, Bedrock-free) =====
#
# 게이트웨이 이미지 라우트가 차단되거나, 폴더 구조/플로우차트처럼
# "사실적 구조"를 그려야 할 때 Stability/Titan 프롬프트 결과보다 훨씬 정확한
# 다이어그램을 로컬에서 만들어준다. matplotlib만 있으면 인터넷/Bedrock 없이 동작.
def _looks_structural(description: str = "", title: str = "", body: str = "") -> bool:
    """Return True iff text contains a literal structural signal — not just
    a generic keyword.

    Spec: media-output-quality (bugfix) — Property 3 / Req 2.2, 3.1.
    Three signals (any one is sufficient):
      1. Path token — at least one identifier-shaped pair separated by `/`.
      2. Arrow chain — `->` / `→` / `⇒` between two non-whitespace tokens.
      3. Markdown table row — a line wrapped in `|` with ≥ 2 cell separators.

    Rationale: the prior keyword-OR matcher caught generic words like
    "프로젝트", "구조", "흐름도", "diagram", "architecture" and silently
    routed visual-intent PPTX/PDF requests to the matplotlib path,
    bypassing real Bedrock image models even when the gateway was healthy.
    """
    import re as _re
    text = "\n".join(s for s in (description, title, body) if s)
    if not text:
        return False
    # Signal 1: path token (e.g. "src/components/foo.js")
    if _re.search(r"[A-Za-z0-9_.\-]+/[A-Za-z0-9_.\-]+", text):
        return True
    # Signal 2: arrow chain (e.g. "A -> B", "A → B", "A ⇒ B")
    if _re.search(r"\S+\s*(->|→|⇒)\s*\S+", text):
        return True
    # Signal 3: markdown table row (≥ 2 pipe-separated cells means ≥ 3 pipes)
    for line in text.splitlines():
        if _re.match(r"^\s*\|.+\|\s*$", line) and line.count("|") >= 3:
            return True
    return False


# ===== Mermaid Diagram Generation =====
# matplotlib보다 훨씬 깔끔한 다이어그램. Claude Sonnet이 mermaid 코드를 정확히 생성하므로
# LLM 호출 1회 + mermaid.ink HTTP GET 1회로 고품질 PNG를 받는다. 실패 시 matplotlib fallback.

_MERMAID_INK_URL = "https://mermaid.ink/img/{encoded}?type=png&bgColor=ffffff"


def _encode_mermaid_for_ink(mermaid_code: str) -> str:
    """mermaid.ink가 받는 형식(base64 URL-safe)으로 인코딩."""
    import base64 as _b64
    raw = mermaid_code.encode("utf-8")
    return _b64.urlsafe_b64encode(raw).decode("ascii").rstrip("=")


async def _llm_generate_mermaid(
    gw,
    model_id: str,
    section_heading: str,
    section_body: str,
    doc_context: str = "",
    style_profile: dict | None = None,
) -> str:
    """LLM에게 섹션 콘텐츠를 분석해서 mermaid 코드 1개를 생성하게 한다.

    프롬프트는 mermaid 문법 + 한국어 라벨 안전 처리 가이드를 포함. 응답은 mermaid 코드만
    추출 (마크다운 fence 제거).

    Active_Template Style_Profile이 주어지고 primaryColor·textColor가 모두 유효한
    `#RRGGBB`이면(요구사항 7.2), (a) 프롬프트에 해당 색상을 classDef에 쓰도록 명시하고
    (b) 추출된 mermaid 코드 앞에 `%%{{init: ...}}%%` 테마 지시문을 주입해 LLM 출력과
    무관하게 테마 색을 강제한다. style_profile이 None/무효이거나 색상 검증이 실패하면
    프롬프트·출력 모두 기존 동작과 바이트 단위로 동일하다. 어떤 경우에도 예외를 올리지
    않으며, 문제 발생 시 주입 없이 폴백한다.
    """
    # Style_Profile primary/text 색 검증 (요구사항 7.2) — 둘 다 유효한 #RRGGBB일 때만 주입.
    # normalize_color는 비문자열/형식 위반에 None을 반환하므로 안전하다. dual-path import는
    # repo 루트 실행(ai_engine.style_profile) / ai_engine 내부 실행(style_profile) 모두 지원.
    _sp_primary: str | None = None
    _sp_text: str | None = None
    if isinstance(style_profile, dict):
        try:
            try:
                from ai_engine.style_profile import normalize_color as _norm_color
            except ImportError:  # pragma: no cover - alt path when run from ai_engine/
                from style_profile import normalize_color as _norm_color  # type: ignore
            _p = _norm_color(style_profile.get("primaryColor"))
            _t = _norm_color(style_profile.get("textColor"))
            if _p and _t:
                _sp_primary, _sp_text = _p, _t
        except Exception:
            # 어떤 import/검증 실패에도 주입 없이 기존 동작으로 폴백 (절대 raise 금지).
            _sp_primary, _sp_text = None, None

    # 프롬프트 색상 지시문 — 유효 색이 있을 때만 비어있지 않다. None이면 빈 문자열이라
    # 프롬프트가 기존과 바이트 단위로 동일하다 (요구사항 7.2, 폴백 격리).
    _color_instruction = ""
    if _sp_primary and _sp_text:
        _color_instruction = (
            f"\n6. classDef 색상은 주 색상(primary) {_sp_primary}와 "
            f"텍스트 색상 {_sp_text}을 사용하세요. 예: "
            f"classDef primary fill:{_sp_primary},color:{_sp_text}"
        )

    prompt = f"""다음 섹션을 분석해 가장 적합한 mermaid 다이어그램 코드를 생성하세요.

섹션 제목: {section_heading or "(없음)"}
섹션 본문:
{(section_body or "")[:2000]}

문서 컨텍스트: {(doc_context or "")[:500]}

요구사항:
1. mermaid 다이어그램 1개만 생성. 다른 설명 절대 금지.
2. 다이어그램 종류는 콘텐츠에 가장 적합한 것을 자동 선택:
   - 폴더 구조/디렉토리: graph TD (위→아래 트리)
   - 시스템 아키텍처/계층: graph TB with subgraph (계층별 그룹)
   - 데이터 흐름/프로세스: flowchart LR (좌→우 화살표) 또는 sequenceDiagram
   - 기술 스택: graph LR with subgraph 또는 classDiagram
   - 상태 전환: stateDiagram-v2
   - 의존성: graph LR
3. 노드 라벨은 한국어로 작성 가능. 라벨에 콜론, 큰따옴표는 피할 것.
   - 좋은 예: A[프론트엔드<br/>React 18]
   - 나쁜 예: A["프론트엔드: React"] (콜론은 mermaid 파서 오류 유발)
4. 노드 6~12개 권장. 너무 적거나 많으면 시각적으로 비효율.
5. classDef로 색상 부여:
   - classDef primary fill:#cfe2f3,stroke:#3c78d8,color:#1e1e1e
   - classDef accent fill:#fff2cc,stroke:#bf9000,color:#1e1e1e
   - classDef data fill:#d9ead3,stroke:#6aa84f,color:#1e1e1e{_color_instruction}

응답 형식 - mermaid 코드만, 마크다운 fence(```) 없이:"""

    try:
        result = await asyncio.wait_for(
            gw.converse(
                model_id,
                [{"role": "user", "content": [{"text": prompt}]}],
                "당신은 정확한 mermaid 다이어그램 코드를 생성하는 전문가입니다.",
            ),
            timeout=45.0,
        )
    except asyncio.TimeoutError:
        return ""
    except Exception as e:
        print(f"[Mermaid LLM] gateway 호출 실패: {e}")
        return ""

    if not result or result.get("decision") != "ALLOW":
        return ""
    content = result.get("output", {}).get("message", {}).get("content", [])
    raw = "\n".join(c.get("text", "") for c in content if "text" in c).strip()
    if not raw:
        return ""

    # 마크다운 fence 제거
    if "```" in raw:
        # ```mermaid ... ``` 패턴 추출
        import re as _re
        m = _re.search(r"```(?:mermaid)?\s*\n(.+?)\n```", raw, _re.DOTALL)
        if m:
            raw = m.group(1).strip()
        else:
            raw = raw.replace("```mermaid", "").replace("```", "").strip()

    # mermaid 키워드로 시작하는지 검증
    valid_starters = (
        "graph", "flowchart", "sequenceDiagram", "classDiagram",
        "stateDiagram", "stateDiagram-v2", "erDiagram", "gantt",
        "pie", "journey", "mindmap", "timeline",
    )
    first_line = raw.split("\n")[0].strip()
    if not any(first_line.startswith(s) for s in valid_starters):
        # LLM이 추가 설명을 붙였을 수 있음 — 첫 mermaid 라인 찾기
        for ln in raw.split("\n"):
            if any(ln.strip().startswith(s) for s in valid_starters):
                idx = raw.index(ln)
                raw = raw[idx:].strip()
                break
        else:
            return ""

    # Style_Profile 테마 변수 주입 (요구사항 7.2) — primary/text 모두 유효할 때만.
    # LLM 출력의 classDef와 무관하게 테마 색을 강제하기 위해, 이미 `%%{init`로
    # 시작하지 않는 경우에만 init 지시문을 prepend 한다. 색이 없으면 raw 그대로
    # 반환하므로 기존 동작과 바이트 단위로 동일하다.
    if _sp_primary and _sp_text and not raw.lstrip().startswith("%%{init"):
        _init = (
            "%%{init: {'theme':'base','themeVariables':{"
            f"'primaryColor':'{_sp_primary}',"
            f"'primaryTextColor':'{_sp_text}',"
            f"'lineColor':'{_sp_primary}',"
            f"'textColor':'{_sp_text}'"
            "}}}%%\n"
        )
        raw = _init + raw

    return raw


async def _render_mermaid_to_png(
    mermaid_code: str,
    project_path: str = "",
    timeout: int = 30,
) -> str:
    """mermaid 코드를 mermaid.ink 공개 API로 PNG 렌더링.

    Returns:
        JSON 문자열 {path, model: "mermaid", width, height, sizeBytes} 성공 시,
        {error, detail} 실패 시.

    실패 시 한 번 자동으로 sanitize 후 재시도 (슬래시/콜론/따옴표 escape).
    """
    import time as _t
    import httpx as _httpx
    import re as _re

    if not mermaid_code or not mermaid_code.strip():
        return json.dumps({"error": "invalid-parameter", "detail": "mermaid_code is required"})

    def _sanitize_mermaid(code: str) -> str:
        """mermaid 라벨에서 파서 오류를 일으키는 문자 escape.

        - `[/src]` 같은 라벨의 시작 슬래시 → `["/src"]` 로 변경
        - 콜론 `:` → middle dot `·` (선행 콜론은 stateDiagram에서 의미 있으니 라벨 안만)
        """
        out_lines = []
        for ln in code.splitlines():
            # `[xxx]`, `(xxx)`, `{xxx}` 라벨 안에 슬래시/콜론이 있으면 따옴표로 감쌈
            def _wrap(m):
                open_b, content, close_b = m.group(1), m.group(2), m.group(3)
                # 이미 따옴표면 그대로
                if content.startswith('"') and content.endswith('"'):
                    return m.group(0)
                if "/" in content or ":" in content:
                    # 큰따옴표는 \"로 escape
                    safe = content.replace('"', '\\"')
                    return f'{open_b}"{safe}"{close_b}'
                return m.group(0)
            # 노드 라벨 패턴: ID[content], ID(content), ID{content}
            ln = _re.sub(r"(\[)([^\[\]]+?)(\])", _wrap, ln)
            ln = _re.sub(r"(\()([^()]+?)(\))", _wrap, ln)
            ln = _re.sub(r"(\{)([^{}]+?)(\})", _wrap, ln)
            out_lines.append(ln)
        return "\n".join(out_lines)

    async def _attempt(code: str):
        encoded = _encode_mermaid_for_ink(code)
        url = _MERMAID_INK_URL.format(encoded=encoded)
        try:
            async with _httpx.AsyncClient(timeout=_httpx.Timeout(float(timeout), connect=10.0)) as client:
                resp = await client.get(url)
            return resp
        except _httpx.TimeoutException:
            return None
        except Exception:
            return None

    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    ts = str(int(_t.time() * 1000))
    filename = f"mermaid-{ts}.png"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    # 1차 시도: 원본 코드
    resp = await _attempt(mermaid_code)
    if resp is None:
        return json.dumps({"error": "mermaid-render-failed", "detail": f"timeout after {timeout}s"})
    if resp.status_code != 200:
        # 2차 시도: sanitize 후 재시도
        sanitized = _sanitize_mermaid(mermaid_code)
        if sanitized != mermaid_code:
            print(f"[Mermaid] 1차 실패 (HTTP {resp.status_code}) → sanitize 후 재시도")
            resp2 = await _attempt(sanitized)
            if resp2 is not None and resp2.status_code == 200:
                resp = resp2
            else:
                detail = f"HTTP {resp.status_code}: {resp.text[:150]}"
                if resp2 is not None:
                    detail += f" | retry HTTP {resp2.status_code}: {resp2.text[:100]}"
                return json.dumps({"error": "mermaid-render-failed", "detail": detail})
        else:
            return json.dumps({
                "error": "mermaid-render-failed",
                "detail": f"HTTP {resp.status_code}: {resp.text[:200]}",
            })

    if not resp.content or len(resp.content) < 200:
        return json.dumps({"error": "mermaid-render-failed", "detail": "empty PNG response"})
    with open(output_path, "wb") as f:
        f.write(resp.content)

    # 이미지 dimensions
    try:
        from PIL import Image as _PIL
        with _PIL.open(output_path) as im:
            aw, ah = im.size
    except Exception:
        aw, ah = 800, 600

    return json.dumps({
        "path": relative_path,
        "model": "mermaid (mermaid.ink)",
        "width": aw,
        "height": ah,
        "sizeBytes": os.path.getsize(output_path),
    })


# ===== HTML Slide Rendering (Genspark/Gamma-class) =====
# Rationale: matplotlib/Mermaid are great for *technical diagrams* but produce
# the "left text + right tiny picture" look that the user explicitly rejected.
# To match Genspark/Gamma, we render a real 1920x1080 HTML/CSS layout in
# Electron's hidden BrowserWindow and capture it to PNG. The PNG is then used
# as a *full-bleed slide background* in PPTX (left=0, top=0, full size) or as
# a max-12cm image in PDF.
#
# Pipeline:
#   1. _llm_pick_slide_layout — Claude picks one of 7 layouts + JSON data
#   2. slide_templates.render_layout(layout, data) — builds HTML
#   3. _render_html_slide_to_png — POSTs to bridge /bridge/render-html-to-png
#   4. caller passes the PNG path as imageFile= to _tool_generate_pptx/pdf

def _find_local_chrome() -> str:
    """로컬 Chrome/Chromium/Edge 실행 파일 경로 탐지. 없으면 "".

    Electron 브리지가 없어도(독립 uvicorn 등) 서버가 직접 HTML→PNG를 렌더하기 위함.
    AE_CHROME_PATH 환경변수로 명시 지정 가능.
    """
    import shutil as _sh
    cands = [
        os.environ.get("AE_CHROME_PATH", "").strip(),
        "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome",
        "/Applications/Chromium.app/Contents/MacOS/Chromium",
        "/Applications/Microsoft Edge.app/Contents/MacOS/Microsoft Edge",
        "/usr/bin/google-chrome", "/usr/bin/google-chrome-stable",
        "/usr/bin/chromium", "/usr/bin/chromium-browser", "/usr/bin/microsoft-edge",
        "/snap/bin/chromium",
    ]
    for c in cands:
        if c and os.path.isfile(c) and os.access(c, os.X_OK):
            return c
    for name in ("google-chrome", "google-chrome-stable", "chromium",
                 "chromium-browser", "chrome", "microsoft-edge"):
        p = _sh.which(name)
        if p:
            return p
    return ""


def _html_renderer_available():
    """HTML→PNG 렌더러 가용성 판정. (available, renderer) 튜플 반환.

    renderer: Electron 브리지가 살아있으면 "bridge", 아니면 로컬 Chrome이
    탐지되면 "local-chrome", 둘 다 없으면 "".
    추가 네트워크 호출 없음 — 기존 두 경로(_call_bridge("status") + _find_local_chrome)만
    재사용한다(generate_pptx 게이트와 동일 비용).
    """
    try:
        if _call_bridge("status", {}):
            return (True, "bridge")
    except Exception:
        pass
    if _find_local_chrome():
        return (True, "local-chrome")
    return (False, "")


def _render_html_via_local_chrome(html, output_path, width=1920, height=1080, timeout=30) -> dict:
    """로컬 Chrome 헤드리스로 HTML을 PNG로 캡쳐. Electron 무관.

    Returns {ok, path, width, height, sizeBytes, renderer} | {ok: False, error}.
    어떤 예외도 raise하지 않는다(호출자 네이티브 폴백 보장).
    """
    chrome = _find_local_chrome()
    if not chrome:
        return {"ok": False, "error": "no local chrome"}
    import tempfile as _tf, subprocess as _sp, uuid as _uuid
    tmp_html = os.path.join(_tf.gettempdir(), f"ae-slide-{_uuid.uuid4().hex}.html")
    try:
        with open(tmp_html, "w", encoding="utf-8") as _f:
            _f.write(html)
        base = [chrome, "--headless=new", "--disable-gpu", "--hide-scrollbars",
                "--no-sandbox", "--force-device-scale-factor=1",
                f"--window-size={int(width)},{int(height)}",
                f"--screenshot={output_path}", f"file://{tmp_html}"]
        try:
            _sp.run(base, capture_output=True, timeout=float(timeout) + 10.0)
        except Exception:
            pass
        if not (os.path.isfile(output_path) and os.path.getsize(output_path) > 1024):
            # 구버전 Chrome — legacy headless 플래그 재시도
            legacy = list(base); legacy[1] = "--headless"
            try:
                _sp.run(legacy, capture_output=True, timeout=float(timeout) + 10.0)
            except Exception:
                pass
        if os.path.isfile(output_path) and os.path.getsize(output_path) > 1024:
            return {"ok": True, "path": output_path, "width": int(width),
                    "height": int(height), "sizeBytes": os.path.getsize(output_path),
                    "renderer": "local-chrome"}
        return {"ok": False, "error": "chrome produced empty/no png"}
    except Exception as e:
        return {"ok": False, "error": f"local chrome failed: {e}"}
    finally:
        try:
            os.remove(tmp_html)
        except OSError:
            pass


async def _render_html_slide_to_png(
    html: str,
    output_path: str,
    width: int = 1920,
    height: int = 1080,
    timeout: int = 30,
) -> dict:
    """Capture an HTML document to PNG via the Electron bridge.

    The bridge endpoint `/bridge/render-html-to-png` is wired to a hidden
    BrowserWindow on the Electron side (see electron/src/ipc-slides-handler.js).
    When Electron is not running (CI / headless tests / `python -m ai_engine`
    started before Electron), the bridge is unreachable and we return an
    {ok: False} envelope so the caller can fall back to mermaid/matplotlib.

    Args:
        html: complete self-contained HTML document. External http(s) URLs
              will be rejected by the bridge — keep all CSS inline.
        output_path: ABSOLUTE path. Bridge writes the PNG here.
        width:  pixel width of the BrowserWindow (default 1920).
        height: pixel height (default 1080 — 16:9).
        timeout: seconds. Bridge has its own internal timeout slightly above this.

    Returns:
        {ok: True, path, width, height, sizeBytes} on success
        {ok: False, error: str}                    on failure
    """
    if not isinstance(html, str) or not html.strip():
        return {"ok": False, "error": "html is required"}
    if not isinstance(output_path, str) or not os.path.isabs(output_path):
        return {"ok": False, "error": "outputPath must be absolute"}

    # Bridge POST is synchronous (httpx) — wrap in run_in_executor so the event
    # loop is free for other awaits during the multi-second capture.
    payload = {
        "html": html,
        "width": int(width),
        "height": int(height),
        "outputPath": output_path,
        "timeoutMs": int(timeout) * 1000,
    }
    loop = asyncio.get_event_loop()
    try:
        result = await loop.run_in_executor(
            None,
            lambda: _call_bridge("render-html-to-png", payload, timeout=float(timeout) + 10.0),
        )
    except Exception as e:
        return {"ok": False, "error": f"bridge call failed: {e}"}

    if not result:
        # Bridge unreachable (Electron not running) — 로컬 Chrome 헤드리스로 직접 렌더.
        local = await loop.run_in_executor(
            None,
            lambda: _render_html_via_local_chrome(html, output_path, int(width), int(height), int(timeout)),
        )
        if local and local.get("ok"):
            return local
        return {"ok": False, "error": "bridge unreachable and local chrome render failed"}
    if not isinstance(result, dict):
        return {"ok": False, "error": "bridge returned non-dict"}
    return result


async def _llm_pick_slide_layout(
    gw,
    model_id: str,
    section_heading: str,
    section_body: str,
    doc_context: str = "",
    role: str = "",
    bullet_count: int = 0,
) -> dict:
    """Ask Claude to pick the best slide layout for this section + extract data.

    Returns:
        {layout: <name>, data: <dict>} on success.
        {layout: "", data: {}}        if no layout fits or LLM call failed.

    The caller (`_force_generate_from_text`) treats an empty layout as "fall
    back to mermaid/matplotlib" — so this function NEVER raises; on any error
    it returns the empty envelope.
    """
    if not gw or not model_id:
        return {"layout": "", "data": {}}
    if not (section_heading or section_body):
        return {"layout": "", "data": {}}

    # Prompt is structured: enumerate layouts with their required data shapes.
    # Mixing English (for the schema) and Korean (for the instruction) keeps
    # Claude focused on the JSON contract while still understanding KR content.
    # 고밀도 힌트(task 3.4): 역할(role)·불릿 수를 근거로 참고 매뉴얼 수준의 고밀도
    # 레이아웃(STEP 카드/2단/KPI/타임라인/상태표/프로세스 흐름)을 적극 선택하도록 유도.
    _bc = int(bullet_count or 0)
    _density_hint = (
        f"\n[밀도 가이드] 이 슬라이드 역할: {role or 'content'} · 불릿 {_bc}개.\n"
        "- 본문(content) 슬라이드는 '빈 텍스트 나열'이 아니라 참고 매뉴얼 수준의 고밀도\n"
        "  레이아웃을 우선 선택한다. 단조로운 단일 컬럼/장문 텍스트는 피한다.\n"
        "- 불릿/항목 4개 이상이면 two_column·feature_grid·status_table 같은 다열 고밀도\n"
        "  레이아웃을, 단계/순서/프로세스 성격이면 timeline·process_flow를 우선 고려한다.\n"
        "- 정량 지표/수치(%, 건수, 금액)가 보이면 kpi_summary, 항목별 상태/진척이면\n"
        "  status_table, 단일 목표의 상세(근거/증빙)면 objective_detail을 고려한다.\n"
    )
    prompt = f"""다음 섹션 콘텐츠를 분석해 가장 적합한 슬라이드 레이아웃 1개를 선택하고
필요한 데이터를 JSON으로 추출하세요.

섹션 제목: {(section_heading or "(없음)")[:200]}
섹션 본문:
{(section_body or "")[:1800]}

문서 컨텍스트: {(doc_context or "")[:300]}
{_density_hint}
다음 11가지 레이아웃 중 가장 적합한 1개 선택 (본문은 고밀도 레이아웃을 우선):

1. cover — 표지/타이틀 슬라이드 (보통 첫 슬라이드 한 장만)
   data: {{"title": str, "subtitle": str (optional), "eyebrow": str (optional)}}

2. section_divider — 섹션 구분 (예: "01 / 프로젝트 개요")
   data: {{"title": str, "section_number": int, "description": str (optional)}}

3. two_column — 좌우 분할 (좌: 설명/정의, 우: 강점/특징)
   data: {{"title": str, "left_content": str (multiline bullets ok),
           "right_content": str (multiline bullets ok), "subtitle": str (optional)}}

4. feature_grid — 3-6개 카드 그리드 (기능, 장점, 기술 스택 등)
   data: {{"title": str, "features": [{{"icon": str, "title": str, "description": str}}, ...],
           "subtitle": str (optional)}}
   icon은 다음 중 하나: check, arrow_right, layers, zap, shield, code, database, cloud, cpu, users, circle

5. timeline — 단계/순서/프로세스 (2-7 steps)
   data: {{"title": str, "steps": [{{"label": str, "title": str, "description": str}}, ...],
           "subtitle": str (optional)}}

6. comparison — 좌우 비교 (Before/After, A/B, 장단점 등)
   data: {{"title": str, "left_label": str, "left_items": [str, ...],
           "right_label": str, "right_items": [str, ...],
           "subtitle": str (optional)}}

7. architecture — 시스템 계층 구조 (Frontend → Backend → DB 등 2-5 레이어)
   data: {{"title": str, "layers": [{{"name": str, "description": str, "items": [str, ...]}}, ...],
           "subtitle": str (optional)}}

8. kpi_summary — 정량 지표/수치 요약 (2-5개 KPI 카드: 값+라벨)
   data: {{"title": str, "metrics": [{{"value": str, "label": str, "sublabel": str (optional), "tone": str (optional: primary|secondary|accent|dark|neutral)}}, ...], "subtitle": str (optional), "eyebrow": str (optional)}}

9. status_table — 항목별 상태/진척 표 (행마다 셀 + 진척률/상태 배지)
   data: {{"title": str, "columns": [str, ...], "rows": [{{"cells": [str, ...], "progress": int (optional 0-100), "status": str (optional), "status_tone": str (optional)}}, ...], "subtitle": str (optional)}}

10. objective_detail — 단일 목표의 상세 (설명 블록 + 증빙/근거)
    data: {{"title": str, "number": str (optional), "subtitle": str (optional), "status": str (optional), "blocks": [{{"heading": str, "items": [str, ...]}}, ...], "evidence": {{"heading": str (optional), "items": [str, ...]}} (optional)}}

11. process_flow — 가로 프로세스 흐름 (3-6 단계 박스가 화살표로 연결)
    data: {{"title": str, "steps": [{{"title": str, "caption": str (optional), "tone": str (optional)}}, ...], "subtitle": str (optional), "note": str (optional)}}

규칙:
- 콘텐츠가 위 어떤 레이아웃에도 적합하지 않으면 layout을 빈 문자열 ""로 반환.
- 본문(content) 슬라이드는 위 [밀도 가이드]를 따라 고밀도 레이아웃을 우선 선택한다.
- title 필드는 섹션 heading을 그대로 또는 더 명확하게 정리해 사용.
- 각 layout이 요구하는 data shape를 정확히 따를 것 (다른 키 추가 금지).
- 모든 텍스트는 한국어 그대로 유지 (영어 번역 금지).
- 출력은 순수 JSON만. 마크다운 fence(```) 없이. 다른 설명 없이.

응답 형식:
{{"layout": "feature_grid", "data": {{"title": "...", "features": [...]}}}}"""

    try:
        result = await asyncio.wait_for(
            gw.converse(
                model_id,
                [{"role": "user", "content": [{"text": prompt}]}],
                "당신은 문서 콘텐츠를 슬라이드 레이아웃으로 매핑하는 전문가입니다. JSON만 반환합니다.",
            ),
            timeout=45.0,
        )
    except asyncio.TimeoutError:
        return {"layout": "", "data": {}}
    except Exception as e:
        print(f"[SlideLayout LLM] gateway 호출 실패: {e}")
        return {"layout": "", "data": {}}

    if not result or result.get("decision") != "ALLOW":
        return {"layout": "", "data": {}}
    content = result.get("output", {}).get("message", {}).get("content", [])
    raw = "\n".join(c.get("text", "") for c in content if "text" in c).strip()
    if not raw:
        return {"layout": "", "data": {}}

    # Strip markdown fences if Claude added them despite the instruction
    if "```" in raw:
        m = re.search(r"```(?:json)?\s*\n?(.+?)\n?```", raw, re.DOTALL)
        if m:
            raw = m.group(1).strip()
        else:
            raw = raw.replace("```json", "").replace("```", "").strip()

    # Find first { ... } block in case there's prose before/after
    if not raw.startswith("{"):
        m = re.search(r"\{.*\}", raw, re.DOTALL)
        if m:
            raw = m.group(0)

    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError as e:
        print(f"[SlideLayout LLM] JSON parse error: {e} | raw: {raw[:200]}")
        return {"layout": "", "data": {}}

    if not isinstance(parsed, dict):
        return {"layout": "", "data": {}}
    layout = (parsed.get("layout") or "").strip()
    data = parsed.get("data") or {}
    if not isinstance(data, dict):
        return {"layout": "", "data": {}}

    # Validate against known layouts — slide_templates registry is the
    # source of truth. If layout name is unknown, treat as "no fit".
    try:
        from ai_engine.slide_templates import LAYOUT_REGISTRY
    except ImportError:
        try:
            from slide_templates import LAYOUT_REGISTRY  # alt path
        except ImportError:
            LAYOUT_REGISTRY = {}
    if layout and layout not in LAYOUT_REGISTRY:
        print(f"[SlideLayout LLM] unknown layout: {layout}")
        return {"layout": "", "data": {}}

    return {"layout": layout, "data": data}


async def _llm_structure_native_diagram(gw, model_id, section_heading, section_body, doc_context=""):
    """섹션을 편집 가능한 다이어그램 스펙으로 구조화. 반환 {type, content} 또는 {}.
    content: flow='A -> B -> C', tree=들여쓰기(2칸=1depth), block=줄바꿈 항목."""
    if not gw or not model_id:
        return {}
    if not (section_heading or section_body):
        return {}
    prompt = (
        "다음 슬라이드 섹션을 PowerPoint에서 편집 가능한 다이어그램로 구조화하세요.\n\n"
        "섹션 제목: " + str(section_heading or "")[:200] + "\n섹션 본문:\n" + str(section_body or "")[:1500] + "\n\n"
        "규칙:\n"
        "- type 하나 선택: flow(순서/흐름/프로세스), tree(구조/계층/조직), cards(특징/장점/항목 설명 2~6개), block(단순 나열).\n"
        "- content 포맷: flow는 '노드1 -> 노드2 -> 노드3', tree는 들여쓰기 2칸=1단계(첫줄=루트), "
        "cards는 한 줄에 '제목: 설명'(결하나 당 1개) 형식 2~6줄, block은 줄바꿈 항목 3~6개.\n"
        "- 노드 라벨은 의미 통하는 짧은 명사구(2~12자). 문장/조사/단어파편 금지.\n"
        "- (선택) note: 해당 섹션에 꾭 필요한 주의·핵심 안내 1문장이 있으면 note 필드에 넣는다(없으면 생략).\n"
        "- 부적합하면 type을 ""로. 순수 JSON만 출력(마크다운 fence 없이).\n\n"
        '응답 예: {"type": "cards", "content": "체계적 구조: 93개 디렉토리\\n효율적 분류: 5개 카테고리\\n확장 가능: 멀티레벨 구조"}'
    )
    try:
        result = await asyncio.wait_for(
            gw.converse(
                model_id,
                [{"role": "user", "content": [{"text": prompt}]}],
                "당신은 슬라이드를 편집 가능한 다이어그램 구조로 변환하는 전문가입니다. JSON만 반환.",
            ),
            timeout=45.0,
        )
    except asyncio.TimeoutError:
        return {}
    except Exception as e:
        print(f"[NativeStruct LLM] gateway 실패: {e}")
        return {}
    if not result or result.get("decision") != "ALLOW":
        return {}
    content = result.get("output", {}).get("message", {}).get("content", [])
    raw = "\n".join(c.get("text", "") for c in content if "text" in c).strip()
    if not raw:
        return {}
    if "```" in raw:
        m = re.search(r"```(?:json)?\s*\n?(.+?)\n?```", raw, re.DOTALL)
        raw = m.group(1).strip() if m else raw.replace("```json", "").replace("```", "").strip()
    if not raw.startswith("{"):
        m = re.search(r"\{.*\}", raw, re.DOTALL)
        if m:
            raw = m.group(0)
    try:
        parsed = json.loads(raw)
    except json.JSONDecodeError:
        return {}
    if not isinstance(parsed, dict):
        return {}
    dtype = (parsed.get("type") or "").strip().lower()
    dcontent = (parsed.get("content") or "").strip()
    if dtype not in ("flow", "tree", "block", "cards") or not dcontent:
        return {}
    _dnote = (parsed.get("note") or "").strip()
    return {"type": dtype, "content": dcontent, "note": _dnote}


def _heuristic_html_layout(heading, body, role="", bullet_count=0):
    """LLM 레이아웃 선정 실패 시 결정론적 폴백 — 불릿을 feature_grid 카드로 매핑.

    빈 슬라이드를 방지하고 본문도 항상 템플릿 스타일 HTML(카드 그리드)로 렌더되게 한다.
    아이콘은 슬라이드 ICONS 키(zap/cloud/shield/database/users/code/layers/check) 중 선택.
    """
    import re as _re2
    lines = [ln.strip().lstrip("-*\u2022 0123456789.").strip()
             for ln in (body or "").splitlines()]
    lines = [l for l in lines if l]
    if not lines:
        lines = [s.strip() for s in _re2.split(r"[.\n\u00b7]", body or "") if s.strip()][:6]
    _kw = [("zap", ("\uc131\ub2a5", "\uc18d\ub3c4", "\ube60\ub978", "\uc2e4\uc2dc\uac04", "\uc989\uc2dc", "fast", "speed")),
           ("cloud", ("\ud074\ub77c\uc6b0\ub4dc", "aws", "\uc11c\ubc84", "\uc778\ud504\ub77c", "m365", "teams", "onedrive", "cloud")),
           ("shield", ("\ubcf4\uc548", "\uad8c\ud55c", "\ud1b5\uc81c", "\uac10\uc0ac", "\uc778\uc99d", "security")),
           ("database", ("\ub370\uc774\ud130", "\uc800\uc7a5", "\ubc31\uc5c5", "db", "database")),
           ("users", ("\uc0ac\uc6a9\uc790", "\uacc4\uc815", "\ud300", "\uba64\ubc84", "user", "team")),
           ("code", ("\ucf54\ub4dc", "\uac1c\ubc1c", "\uc124\uce58", "\ube4c\ub4dc", "\uc2a4\ud06c\ub9bd\ud2b8", "code", "install")),
           ("layers", ("\uad6c\uc870", "\uacc4\uce35", "\uc544\ud0a4\ud14d\ucc98", "\uad6c\uc131", "layer"))]
    def _ico(t):
        tl = (t or "").lower()
        for ic, ks in _kw:
            if any(k in tl for k in ks):
                return ic
        return "check"
    feats = []
    for ln in lines[:6]:
        title, desc = ln, ""
        for sep in (" \u2014 ", " - ", " | ", ": ", "\uff1a", "\u2014", "|"):
            if sep in ln:
                a, b = ln.split(sep, 1)
                title, desc = a.strip(), b.strip()
                break
        feats.append({"icon": _ico(title + " " + desc),
                      "title": title[:40], "description": desc[:90]})
    if not feats:
        return "", {}
    _title = (heading or "\uc8fc\uc694 \ub0b4\uc6a9")[:60]
    _n = len(feats)
    # 고밀도 폴백(task 3.4): 단계/순서/프로세스 성격 → process_flow, 4+ 항목 → two_column,
    # 그 외 → feature_grid. 참고 매뉴얼 수준의 본문 밀도를 결정론적으로 보장한다.
    _sig = ((heading or "") + " " + (body or "")).lower()
    _step_kw = ("단계", "순서", "프로세스", "절차", "흐름", "step", "process", "phase")
    _has_step = any(k in _sig for k in _step_kw)
    if _has_step and 3 <= _n <= 6:
        _steps = [{"title": f["title"], "caption": f.get("description", "")} for f in feats]
        return "process_flow", {"title": _title, "steps": _steps}
    if _n >= 4:
        _half = (_n + 1) // 2
        def _line(f):
            return f["title"] + (": " + f["description"] if f.get("description") else "")
        _left = "\n".join(_line(f) for f in feats[:_half])
        _right = "\n".join(_line(f) for f in feats[_half:])
        return "two_column", {"title": _title, "left_content": _left, "right_content": _right}
    return "feature_grid", {"title": _title, "features": feats}


def _cards_fallback_from_bullets(bullets):
    """불릿 리스트 → 편집 가능 네이티브 카드 다이어그램 spec(또는 None).

    HTML 고밀도 경로가 비활성이거나 렌더에 실패했을 때 본문 슬라이드가 '휑한 불릿'
    으로 떨어지지 않도록, 4+ 항목은 twocol, 2-3 항목은 cards 로 매핑한다(task 3.4).
    HTML 활성·성공 경로에서는 호출되지 않으므로 HTML 풀블리드 우선순위를 보존한다.
    """
    import re as _re_cf2
    _items = []
    for _b in (bullets or []):
        _t = str(_b).strip()
        if not _t:
            continue
        _t = _re_cf2.sub(r"`([^`]+)`", r"\1", _t)
        _t = _t.replace("**", "").replace("__", "")
        _t = _re_cf2.sub(r"^\s*(?:[-*•]\s+|#{1,6}\s+)", "", _t)
        if _t:
            _items.append(_t)
    if len(_items) == 1 and len(_items[0]) > 40:
        _segs = _re_cf2.split(r"(?:\n|·|•|;|\. |。)", _items[0])
        _segs = [s.strip(" .") for s in _segs if len(s.strip()) > 3]
        if len(_segs) >= 2:
            _items = _segs
    _items = _items[:6]
    if len(_items) >= 4:
        return {"type": "twocol", "content": "\n".join(_items)}
    if len(_items) >= 2:
        return {"type": "cards", "content": "\n".join(_items)}
    return None


async def _generate_html_slide_for_section(
    gw,
    model_id: str,
    section_heading: str,
    section_body: str,
    doc_context: str,
    project_path: str,
    style_profile=None,
    hero_image: str = "",
    render_info: Optional[dict] = None,
) -> str:
    """High-level helper: pick layout → render HTML → capture PNG.

    Returns the relative path (`.generated/slide-...png`) on success, or
    empty string on any failure (caller falls back to mermaid/matplotlib).
    """
    # 고밀도 힌트(task 3.4): 불릿 수 + 역할을 추론해 레이아웃 피커에 전달한다.
    _bn = len([ln for ln in (section_body or "").splitlines() if ln.strip()])
    try:
        _role_hint = _classify_slide_role(
            {"title": section_heading,
             "bullets": [l for l in (section_body or "").splitlines() if l.strip()]},
            False, doc_context)
    except Exception:
        _role_hint = "content"
    # 1. Pick layout (역할·불릿 수 힌트로 고밀도 레이아웃을 적극 선택)
    pick = await _llm_pick_slide_layout(
        gw, model_id, section_heading, section_body, doc_context,
        role=_role_hint, bullet_count=_bn,
    )
    layout = pick.get("layout") or ""
    data = pick.get("data") or {}
    if not layout:
        # LLM 실패 → 결정론적 고밀도 폴백(process_flow / two_column / feature_grid)
        layout, data = _heuristic_html_layout(
            section_heading, section_body, role=_role_hint, bullet_count=_bn)
    if not layout:
        return ""
    # Active_Template Style_Profile → HTML 디자인 토큰(색/폰트) 적용
    # (Genspark급 레이아웃 + 템플릿 색상 통합). 무프로파일이면 기본 디자인 유지.
    if style_profile:
        try:
            try:
                from ai_engine.slide_templates import design_tokens_for_profile as _dtfp
            except ImportError:
                from slide_templates import design_tokens_for_profile as _dtfp
            data["design"] = _dtfp(style_profile)
        except Exception:
            pass

    # 1.5 HTML+Vertex 합성(task 3.5) — hero_image가 주어지면 레이아웃의 선택적
    #     이미지 슬롯(cover=heroImage, two_column/objective_detail=image)에 주입해
    #     'HTML 레이아웃 품질 + Vertex 이미지'를 단일 PNG로 통합한다. 슬롯이 없는
    #     레이아웃은 합성하지 않고 render_info로 신호 → caller가 on-slide 레이어링으로
    #     폴백(이미지 폐기 금지). hero_image 미주입 시 data 불변 → 바이트 호환.
    _composited = False
    if hero_image:
        try:
            if layout == "cover":
                data["heroImage"] = hero_image
                _composited = True
            elif layout in ("two_column", "objective_detail"):
                data["image"] = hero_image
                _composited = True
        except Exception:
            _composited = False
    if isinstance(render_info, dict):
        render_info["layout"] = layout
        render_info["composited"] = _composited

    # 2. Render HTML
    try:
        try:
            from ai_engine.slide_templates import render_layout
        except ImportError:
            from slide_templates import render_layout  # alt path when running from ai_engine/
    except ImportError as e:
        print(f"[HtmlSlide] slide_templates import failed: {e}")
        return ""
    html = render_layout(layout, data)
    if not html:
        return ""

    # 3. Capture PNG via Electron bridge
    import time as _t
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    ts = str(int(_t.time() * 1000))
    safe_layout = re.sub(r"[^a-zA-Z0-9_]", "", layout)[:20] or "slide"
    filename = f"slide-{safe_layout}-{ts}.png"
    abs_out = os.path.join(gen_dir, filename)
    relative = f".generated/{filename}"

    res = await _render_html_slide_to_png(html, abs_out, width=1920, height=1080, timeout=30)
    if not (res and res.get("ok")):
        err = (res or {}).get("error", "unknown")
        print(f"[HtmlSlide] capture failed (layout={layout}): {err}")
        return ""
    if not os.path.isfile(abs_out) or os.path.getsize(abs_out) < 1000:
        print(f"[HtmlSlide] capture wrote tiny file: {abs_out}")
        return ""
    print(f"[HtmlSlide] section[{section_heading[:40]}] layout={layout} → {relative}")
    return relative


def _normalize_palette(palette):
    """팔레트 입력을 유효 `#RRGGBB` 2색 이상 리스트로 정규화한다 (요구사항 7.3/7.6).

    - 입력이 list/tuple이 아니거나 비어있으면 None.
    - 각 항목을 normalize_color로 검증(대문자 #RRGGBB). 무효 토큰은 제외(요구사항 7.6).
    - 중복은 제거하고 첫 등장 순서를 유지 → primary(첫 유효 색)가 palette[0]로 남는다.
    - 유효 색이 2색 미만이면 None을 반환해 호출자가 기존 기본 색상으로 폴백한다(요구사항 7.5).
    - 어떤 입력에도 예외를 올리지 않는다 (dual-path import 포함).
    """
    if not isinstance(palette, (list, tuple)) or not palette:
        return None
    try:
        try:
            from ai_engine.style_profile import normalize_color as _norm_color
        except ImportError:  # pragma: no cover - alt path when run from ai_engine/
            from style_profile import normalize_color as _norm_color  # type: ignore
    except Exception:
        return None
    out = []
    for c in palette:
        try:
            nc = _norm_color(c)
        except Exception:
            nc = None
        if nc and nc not in out:
            out.append(nc)
    if len(out) < 2:
        return None
    return out


def _build_palette(profile):
    """Style_Profile dict → [primaryColor, secondaryColor, accentColor]에서 유효 `#RRGGBB`만
    (primary 우선, 2색 이상) 추출. 2색 미만이거나 dict가 아니면 None(기본 팔레트 폴백).

    요구사항 7.3(주 색상 첫 항목, 2색 이상) / 7.6(무효 토큰 제외).
    """
    if not isinstance(profile, dict):
        return None
    return _normalize_palette([
        profile.get("primaryColor"),
        profile.get("secondaryColor"),
        profile.get("accentColor"),
    ])


def _hex_to_face_edge(hex_color):
    """팔레트 색(`#RRGGBB`) → (face, edge) 쌍. edge=원색, face=흰색 75% 혼합한 밝은 톤.

    박스 채움(face)을 밝게 유지해 어두운 텍스트(#1e1e1e)가 읽히도록 한다. 무효 입력이면 None.
    """
    try:
        h = hex_color.lstrip("#")
        r = int(h[0:2], 16)
        g = int(h[2:4], 16)
        b = int(h[4:6], 16)
        fr = int(r + (255 - r) * 0.75)
        fg = int(g + (255 - g) * 0.75)
        fb = int(b + (255 - b) * 0.75)
        return (f"#{fr:02x}{fg:02x}{fb:02x}", hex_color)
    except Exception:
        return None


def _palette_face_edges(palette):
    """정규화된 팔레트 → [(face, edge), ...] 리스트. 유효 쌍이 2개 미만이면 None.

    render_* 함수가 박스/카드 색을 팔레트에서 순환 적용할 때 사용한다. palette가 None이면
    None을 반환해 호출자가 기존 하드코딩 색상(바이트 동일)을 쓰도록 한다 (요구사항 7.5).
    """
    if not palette:
        return None
    fe = []
    for c in palette:
        pair = _hex_to_face_edge(c)
        if pair:
            fe.append(pair)
    return fe if len(fe) >= 2 else None


def _build_fullbleed_vertex_prompt(role, title, bullets, style_profile):
    """풀블리드 Vertex 배경 프롬프트 빌더 (순수 결정 함수, task 4.1).

    풀블리드 대상 role(cover/section/visual)별로 서로 다른 프롬프트 본문을 산출하고
    `(prompt, negative_prompt)`를 반환한다. LLM/네트워크 호출 없이 순수 문자열만 조립한다.

    특성(설계 Vertex 프롬프트 빌더 / VertexPrompt / Properties 10,11,12):
    - 역할별 구별성(R3.3): cover/section/visual이 서로 다른 프롬프트 본문을 낸다.
    - 16:9 명시(R3.1): 프롬프트 문말에 "16:9"를 포함한다.
    - no-text negative(R3.2): negative_prompt는 항상 len >= 1이며
      text/words/letters/watermark 억제 용어를 포함한다.
    - 팔레트 결정성(R3.4): `_build_palette(style_profile)`의 primary/secondary 색을
      **고정 키 순서**(primary=palette[0], secondary=palette[1])로 색상 표현에 삽입한다.
      팔레트가 None이면 결정론적 기본 색상 표현으로 폴백한다.
    - 바이트 결정성(R3.6): 동일 입력 → 동일 출력. 난수/타임스탬프/dict 순회에 비의존.
    - 게이트웨이 제약(R3.7): 직접 LLM 호출을 수행하지 않는다.

    어떤 입력에도 raise하지 않는다(방어적 정규화).
    """
    # role 정규화 — 풀블리드 대상(cover/section/visual) 외/모호 입력은 결정론적으로
    # "visual"로 확정한다(풀블리드 히어로 비주얼이 가장 안전한 기본값).
    _role = role if role in ("cover", "section", "visual") else "visual"

    # title/bullets 정규화 (raise 없음, 결정론적)
    _title = str(title or "").strip()[:120]
    _bul = []
    try:
        for _b in (bullets or []):
            _t = str(_b).strip()
            if _t:
                _bul.append(_t)
    except Exception:
        _bul = []
    _topic = (" ".join(_bul))[:500].strip()
    _theme = (_title + (" " + _topic if _topic else "")).strip()

    # 팔레트 색상 표현 — `_build_palette`는 정렬된 리스트([primary, secondary, ...]) 또는
    # None을 반환한다(dict 순회 없음 → 결정론적). 고정 키 순서로 primary/secondary만 사용.
    _pal = _build_palette(style_profile)
    if _pal and len(_pal) >= 2:
        _color_expr = f"palette anchored on {_pal[0]} with {_pal[1]} accent"
    elif _pal and len(_pal) >= 1:
        _color_expr = f"palette anchored on {_pal[0]}"
    else:
        _color_expr = "professional deep navy and blue palette with a single warm accent"

    # 역할별 프롬프트 본문 — 세 역할이 서로 다른 선두 문장을 사용한다(R3.3 구별성 보장).
    if _role == "cover":
        _prompt = (
            f'A commercial-grade hero title background for the cover of a premium corporate '
            f'presentation. Theme: "{_theme}". Style: cinematic depth, dramatic yet elegant '
            f'lighting, expansive negative space in the upper-left for an overlaid title, '
            f'refined executive aesthetic. {_color_expr}. Balanced composition, 16:9.'
        )
    elif _role == "section":
        _prompt = (
            f'A chapter divider ambient background introducing a new section of a corporate '
            f'presentation. Theme: "{_theme}". Style: calm minimalist atmosphere, soft '
            f'gradients and gentle abstract texture, generous empty space for a section '
            f'heading, understated professional mood. {_color_expr}. Balanced composition, 16:9.'
        )
    else:  # visual
        _prompt = (
            f'An editorial photographic hero visual for a corporate presentation slide. '
            f'Theme: "{_theme}". Style: premium professional photography, natural soft '
            f'lighting, shallow depth of field, ample negative space for overlaid text, '
            f'refined corporate aesthetic. {_color_expr}. Balanced composition, 16:9.'
        )

    # no-text negative prompt — 항상 비어 있지 않으며 text/words/letters/watermark 포함(R3.2).
    _negative = (
        "text, words, letters, captions, typography, watermark, fake logo, brand name, "
        "emoji, charts, diagrams, distorted text, unreadable artifacts, childish clipart"
    )
    return _prompt, _negative


async def _tool_generate_native_diagram(
    diagram_type: str,
    title: str,
    content: str,
    project_path: str = "",
    size: tuple = (1600, 1100),
    palette: list | None = None,
) -> str:
    """matplotlib로 폴더 구조/플로우차트를 PNG 직접 생성. Bedrock 무관.

    diagram_type:
        "tree"  → 들여쓰기 깊이 기반 폴더 트리 (각 라인을 박스로)
        "flow"  → "->", "→" 분리 후 좌→우 화살표 흐름
        "block" → 줄바꿈 분리 항목 → 세로 박스 + 화살표

    palette: Active_Template Style_Profile에서 파생된 색상 팔레트(요구사항 7.3).
        유효한 `#RRGGBB` 문자열 2색 이상의 리스트이면 박스/강조 색상에 적용하며
        primary가 첫 항목(palette[0])이 된다. None이거나 유효 색이 2색 미만이면
        기존 하드코딩 기본 색상을 그대로 사용해 출력이 바이트 단위로 동일하다
        (요구사항 7.5/7.6). 어떤 입력에도 예외를 올리지 않는다.

    Returns JSON {path, model: "matplotlib (native)", width, height} or {error}.
    """
    import time as _t

    # 팔레트 정규화 — 유효 #RRGGBB 2색 이상이면 list, 아니면 None(기본 색상 폴백).
    # _normalize_palette는 어떤 입력에도 raise하지 않는다 (요구사항 7.3/7.6).
    _palette = _normalize_palette(palette)

    diagram_type = (diagram_type or "tree").lower()
    title = (title or "Diagram").strip()
    content = (content or "").strip()
    if not content:
        return json.dumps({"error": "invalid-parameter", "detail": "content is required"})

    # matplotlib는 try/except — 미설치 시에도 호출자(force-generate)는
    # 텍스트 전용 PDF/PPTX/DOCX 경로로 fallback 가능
    try:
        import matplotlib
        matplotlib.use("Agg")  # headless backend — 어떤 OS에서도 안전
        import matplotlib.pyplot as plt
        from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
    except ImportError as e:
        return json.dumps({
            "error": "missing-dep",
            "lib": "matplotlib",
            "detail": str(e)[:200],
            "hint": "pip install matplotlib",
        })

    # 한글 폰트 — 시스템에 있으면 사용, 없으면 DejaVu Sans (□ 처리)
    # _apply_korean_font_for_matplotlib()이 _resolve_korean_font_path()로 시스템
    # 폰트를 자동 탐지 + addfont로 등록까지 처리한다 (fc-cache 없는 환경 대응).
    _apply_korean_font_for_matplotlib()

    # 출력 경로 — _tool_generate_image와 같은 _local_root/.generated/ 사용
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    ts = str(int(_t.time() * 1000))
    filename = f"native-{diagram_type}-{ts}.png"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    # 캔버스 — DPI 100 기준 → (1024,768) = (10.24, 7.68) inch
    w_in = max(4, size[0] / 100.0)
    h_in = max(3, size[1] / 100.0)

    try:
        fig, ax = plt.subplots(figsize=(w_in, h_in), dpi=100)
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)
        ax.axis("off")
        # 제목
        ax.text(50, 97, title[:80], ha="center", va="top",
                fontsize=14, fontweight="bold", color="#1e1e1e")

        if diagram_type == "tree":
            _render_tree(ax, content, palette=_palette)
        elif diagram_type == "flow":
            _render_flow(ax, content, palette=_palette)
        elif diagram_type == "architecture":
            _render_architecture(ax, content, palette=_palette)
        elif diagram_type == "stack":
            _render_stack(ax, content, palette=_palette)
        elif diagram_type == "block":
            _render_block(ax, content, palette=_palette)
        else:
            _render_block(ax, content, palette=_palette)

        plt.tight_layout()
        fig.savefig(output_path, dpi=100, bbox_inches="tight",
                    facecolor="white", edgecolor="none")
        plt.close(fig)
    except Exception as e:
        try:
            plt.close("all")
        except Exception:
            pass
        return json.dumps({"error": "diagram-generation-failed", "detail": str(e)[:200]})

    if not os.path.isfile(output_path):
        return json.dumps({"error": "diagram-generation-failed", "detail": "no output file"})

    try:
        from PIL import Image as _PIL
        with _PIL.open(output_path) as im:
            aw, ah = im.size
    except Exception:
        aw, ah = size[0], size[1]

    return json.dumps({
        "path": relative_path,
        "model": "matplotlib (native)",
        "width": aw,
        "height": ah,
        "sizeBytes": os.path.getsize(output_path),
    })


def _render_tree(ax, content: str, palette=None):
    """들여쓰기 기반 폴더 트리 — 부모-자식 가지선(│ ├ └)을 직접 라인으로 그려
    "그림 다이어그램"으로 보이게 만든다. 각 노드는 폴더(하늘색)/파일(회색)로 구분.

    palette가 주어지면(정규화된 `#RRGGBB` 2색 이상) 폴더 박스는 palette[0](primary),
    파일 박스는 palette[1]에서 파생한 색을 사용한다. None이면 기존 하드코딩 색상을 유지해
    출력이 바이트 단위로 동일하다 (요구사항 7.3/7.5).

    파싱:
    - 입력 라인의 들여쓰기 깊이(tab/2-space/4-space)를 자동 측정
    - "│ ├── └──" 같은 트리 문자가 이미 있으면 제거하고 깊이만 사용
    - 마지막 문자가 "/" 또는 "\\"면 폴더로 표시
    """
    from matplotlib.patches import FancyBboxPatch
    from matplotlib.lines import Line2D
    lines_raw = [ln for ln in content.splitlines() if ln.strip()]
    if not lines_raw:
        return
    lines_raw = lines_raw[:40]  # 너무 많으면 cap

    def _depth(s: str) -> int:
        """선행 공백/탭만으로 깊이 측정 — 트리 문자(│├└─)는 nesting을 의미하지 않음."""
        i = 0
        d = 0
        # tab는 4 space로 환산
        while i < len(s):
            if s[i] == "\t":
                d += 1
                i += 1
            elif s.startswith("    ", i):
                d += 1
                i += 4
            elif s.startswith("  ", i):
                # 2-space 들여쓰기 — 0.5 depth (보수적으로 1로 처리)
                d += 1
                i += 2
            elif s[i] in "│ ":
                # 트리 문자 또는 단일 공백은 깊이로 계산하지 않음
                i += 1
            else:
                break
        return d

    # 노드 추출 + 정규화
    nodes = []
    for raw in lines_raw:
        d = _depth(raw)
        # 트리 그래픽 문자 + 선행 공백 제거
        text = raw
        for ch in ("│", "├", "└", "─", "—", "├──", "└──"):
            text = text.replace(ch, "")
        text = text.strip().lstrip("-*•").strip()
        if not text:
            continue
        is_dir = text.endswith("/") or text.endswith("\\")
        nodes.append({"depth": d, "text": text.rstrip("/\\"), "is_dir": is_dir})

    if not nodes:
        return

    n = len(nodes)
    # 캔버스: y는 90(상단)~8(하단) 사용, 노드 높이 자동 산정
    top, bottom = 88, 8
    avail = top - bottom
    box_h = max(2.2, min(4.0, avail / max(n, 1) - 0.5))
    y_step = avail / max(n, 1)

    # 각 깊이의 x 시작점(왼쪽 가장자리). 트리는 깊이당 6 단위씩 들여씀.
    indent_unit = 5
    base_x = 8

    # 각 노드의 위치 계산
    positions = []  # (x, y_top, y_center, y_bottom, depth)
    for i, nd in enumerate(nodes):
        x = base_x + nd["depth"] * indent_unit
        y_t = top - i * y_step
        y_b = y_t - box_h
        positions.append({
            "x": x, "y_top": y_t, "y_bot": y_b, "y_center": (y_t + y_b) / 2,
            "depth": nd["depth"], "is_dir": nd["is_dir"], "text": nd["text"],
        })

    # 가지선(branch line) — 각 노드는 자기보다 깊이가 작은 가장 가까운 위쪽 노드(부모)로
    # 세로선 + 수평 stub을 긋는다. 진짜 트리 구조 시각화의 핵심.
    branch_color = "#888888"
    for i, p in enumerate(positions):
        if p["depth"] == 0:
            continue
        # 부모 찾기: 위쪽으로 올라가며 depth가 더 작은 첫 노드
        parent = None
        for j in range(i - 1, -1, -1):
            if positions[j]["depth"] < p["depth"]:
                parent = positions[j]
                break
        if not parent:
            continue
        # 세로선: 부모의 (x_branch, y_bot) → 자식의 y_center
        x_branch = parent["x"] + 1.5  # 부모 박스 왼쪽 안쪽
        # 수직선
        ax.add_line(Line2D(
            [x_branch, x_branch],
            [parent["y_bot"], p["y_center"]],
            color=branch_color, linewidth=0.8,
        ))
        # 수평 stub: x_branch → 자식 박스 왼쪽
        ax.add_line(Line2D(
            [x_branch, p["x"]],
            [p["y_center"], p["y_center"]],
            color=branch_color, linewidth=0.8,
        ))

    # 박스 + 텍스트
    # palette가 있으면 폴더=primary(palette[0]), 파일=palette[1] 파생 색 적용 (요구사항 7.3).
    _fe = _palette_face_edges(palette)
    _dir_face, _dir_edge = (_fe[0] if _fe else ("#cfe2f3", "#3c78d8"))
    _file_face, _file_edge = (_fe[1] if _fe else ("#f3f3f3", "#999999"))
    for p in positions:
        if p["is_dir"]:
            face = _dir_face  # 폴더
            edge = _dir_edge
            # 폴더는 끝에 / 표시 (이모지는 한글 폰트에 없어서 깨지므로 사용 안 함)
            label_text = "[D] " + p["text"][:54] + "/"
        else:
            face = _file_face  # 파일
            edge = _file_edge
            label_text = p["text"][:60]
        # 박스 너비 — 깊이에 따라 줄어들지만 최소 30
        box_w = max(30, 90 - p["x"])
        ax.add_patch(FancyBboxPatch(
            (p["x"], p["y_bot"]), box_w, box_h,
            boxstyle="round,pad=0.18,rounding_size=0.4",
            facecolor=face, edgecolor=edge, linewidth=1.0,
        ))
        ax.text(p["x"] + 1.2, p["y_center"], label_text,
                va="center", ha="left",
                fontsize=8.5, color="#1e1e1e", weight="bold" if p["is_dir"] else "normal")


def _render_flow(ax, content: str, palette=None):
    """A -> B -> C 형태 좌→우 흐름. 너무 길면 자동 줄바꿈.

    palette가 주어지면 박스/화살표 색을 palette[0](primary)에서 파생해 적용한다.
    None이면 기존 하드코딩 색상을 유지한다 (요구사항 7.3/7.5).
    """
    from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
    # palette 색 적용 — primary(palette[0]) 박스 + 화살표. None이면 기존 색(바이트 동일).
    _fe = _palette_face_edges(palette)
    _box_face, _box_edge = (_fe[0] if _fe else ("#cfe2f3", "#3c78d8"))
    _arrow_color = _box_edge if _fe else "#3c78d8"
    # 분리: -> 또는 → 또는 줄바꿈
    raw = content.replace("→", "->").replace("⇒", "->")
    parts = []
    for chunk in raw.split("\n"):
        for p in chunk.split("->"):
            t = p.strip()
            if t:
                parts.append(t[:40])
    if not parts:
        return
    parts = parts[:8]  # 8개로 cap

    n = len(parts)
    box_w = min(18, 90 / max(n, 1) - 2)
    box_h = 8
    gap = (90 - n * box_w) / max(n - 1, 1) if n > 1 else 0
    y = 50 - box_h / 2

    centers = []
    for i, txt in enumerate(parts):
        x = 5 + i * (box_w + gap)
        ax.add_patch(FancyBboxPatch(
            (x, y), box_w, box_h,
            boxstyle="round,pad=0.3,rounding_size=0.7",
            facecolor=_box_face, edgecolor=_box_edge, linewidth=1.5,
        ))
        ax.text(x + box_w / 2, y + box_h / 2, txt, ha="center", va="center",
                fontsize=10, color="#1e1e1e", wrap=True)
        centers.append((x + box_w, y + box_h / 2, x))
    for i in range(n - 1):
        x_end = centers[i][0]
        x_start_next = centers[i + 1][2]
        y_mid = centers[i][1]
        ax.add_patch(FancyArrowPatch(
            (x_end, y_mid), (x_start_next, y_mid),
            arrowstyle="->", mutation_scale=15,
            color=_arrow_color, linewidth=1.5,
        ))


def _render_block(ax, content: str, palette=None):
    """줄바꿈 분리 항목 → 세로 박스 + 아래 화살표.

    palette가 주어지면 박스/화살표 색을 palette[0](primary)에서 파생해 적용한다.
    None이면 기존 하드코딩 색상을 유지한다 (요구사항 7.3/7.5).
    """
    from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
    _fe = _palette_face_edges(palette)
    _box_face, _box_edge = (_fe[0] if _fe else ("#fff2cc", "#bf9000"))
    _arrow_color = _box_edge if _fe else "#bf9000"
    lines = [ln.strip().lstrip("-*•").strip() for ln in content.splitlines()]
    lines = [ln for ln in lines if ln][:10]
    if not lines:
        return
    n = len(lines)
    box_w = 60
    box_h = max(4, 80 / n - 1.5)
    gap = max(1.5, (80 - n * box_h) / max(n, 1))
    x = 20
    y_top = 90

    centers_y = []
    for i, txt in enumerate(lines):
        y = y_top - i * (box_h + gap) - box_h
        ax.add_patch(FancyBboxPatch(
            (x, y), box_w, box_h,
            boxstyle="round,pad=0.3,rounding_size=0.5",
            facecolor=_box_face, edgecolor=_box_edge, linewidth=1.3,
        ))
        ax.text(x + box_w / 2, y + box_h / 2, txt[:80], ha="center", va="center",
                fontsize=10, color="#1e1e1e")
        centers_y.append((x + box_w / 2, y + box_h, y))
    for i in range(n - 1):
        x_c = centers_y[i][0]
        y_top_arrow = centers_y[i][2]
        y_bot_arrow = centers_y[i + 1][1]
        ax.add_patch(FancyArrowPatch(
            (x_c, y_top_arrow), (x_c, y_bot_arrow),
            arrowstyle="->", mutation_scale=15,
            color=_arrow_color, linewidth=1.5,
        ))


def _render_architecture(ax, content: str, palette=None):
    """3-tier / 마이크로서비스 레이어 다이어그램.

    각 레이어를 가로 박스로 그리고, 사이를 위↔아래 양방향 화살표로 연결.
    파싱: 줄당 한 레이어. "Frontend: React, Redux" 형태면 좌측 라벨 + 우측 component 박스들.

    palette가 주어지면 레이어 박스 색을 palette에서 순환 적용(palette[0]=primary 첫 레이어).
    None이면 기존 layer_colors 하드코딩을 유지한다 (요구사항 7.3/7.5).
    """
    from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
    lines = [ln.strip().lstrip("-*•").strip() for ln in content.splitlines() if ln.strip()]
    lines = lines[:6]  # 최대 6개 레이어
    if not lines:
        return
    n = len(lines)
    layer_h = max(8, 75 / n - 1.5)
    gap = 2.5
    layer_w = 84
    x = 8
    y_top = 90

    layer_colors = [
        ("#cfe2f3", "#3c78d8"),  # 하늘 — 프론트엔드
        ("#d9ead3", "#6aa84f"),  # 연두 — API/앱
        ("#fff2cc", "#bf9000"),  # 노랑 — 비즈니스 로직
        ("#f4cccc", "#cc0000"),  # 분홍 — 데이터
        ("#ead1dc", "#a64d79"),  # 라일락 — 인프라
        ("#d9d2e9", "#674ea7"),  # 보라 — 외부
    ]
    # palette가 있으면 레이어 색을 palette 파생 (face,edge) 쌍으로 교체 (요구사항 7.3).
    _fe = _palette_face_edges(palette)
    if _fe:
        layer_colors = _fe

    layer_centers = []
    for i, txt in enumerate(lines):
        face, edge = layer_colors[i % len(layer_colors)]
        y = y_top - i * (layer_h + gap) - layer_h
        ax.add_patch(FancyBboxPatch(
            (x, y), layer_w, layer_h,
            boxstyle="round,pad=0.4,rounding_size=0.6",
            facecolor=face, edgecolor=edge, linewidth=1.5,
        ))
        # 레이어 텍스트 — heading: components 형태면 분리해서 좌측 굵게 + 우측 일반
        if ":" in txt:
            head, comp = txt.split(":", 1)
            ax.text(x + 1.5, y + layer_h / 2, head.strip()[:30],
                    ha="left", va="center", fontsize=11, color="#1e1e1e", weight="bold")
            ax.text(x + 28, y + layer_h / 2, comp.strip()[:80],
                    ha="left", va="center", fontsize=9.5, color="#1e1e1e")
        else:
            ax.text(x + layer_w / 2, y + layer_h / 2, txt[:90],
                    ha="center", va="center", fontsize=10.5, color="#1e1e1e", weight="bold")
        layer_centers.append((x + layer_w / 2, y + layer_h, y))

    # 레이어 간 양방향 화살표 (요청/응답 의미)
    for i in range(n - 1):
        x_c = layer_centers[i][0]
        y_top_arrow = layer_centers[i][2]
        y_bot_arrow = layer_centers[i + 1][1]
        ax.add_patch(FancyArrowPatch(
            (x_c - 3, y_top_arrow), (x_c - 3, y_bot_arrow),
            arrowstyle="->", mutation_scale=14,
            color="#666", linewidth=1.4,
        ))
        ax.add_patch(FancyArrowPatch(
            (x_c + 3, y_bot_arrow), (x_c + 3, y_top_arrow),
            arrowstyle="->", mutation_scale=14,
            color="#999", linewidth=1.0, linestyle="--",
        ))


def _render_stack(ax, content: str, palette=None):
    """기술 스택 그리드 — 카테고리별 카드. 각 줄을 카드 1개로 매핑.

    파싱:
    - "프론트엔드: React 18, TypeScript 5" → 카테고리 헤더 + 항목들
    - "| 계층 | 기술 | 버전 |" 같은 마크다운 표면 표 우선 파싱

    palette가 주어지면 표 헤더 색과 카드 색을 palette 파생 색으로 교체한다(palette[0]=primary).
    None이면 기존 하드코딩 색상을 유지한다 (요구사항 7.3/7.5).
    """
    from matplotlib.patches import FancyBboxPatch
    # palette 파생 (face,edge) 쌍 — None이면 기존 하드코딩 색상 사용(바이트 동일).
    _fe = _palette_face_edges(palette)
    # 표 헤더 색 — palette가 있으면 primary edge, 없으면 기존 파랑.
    _hdr_face, _hdr_edge = (_fe[0] if _fe else ("#3c78d8", "#1c4587"))
    # 마크다운 표가 있으면 표 우선
    lines_raw = content.splitlines()
    rows = []
    in_table = False
    for ln in lines_raw:
        s = ln.strip()
        if s.startswith("|") and s.endswith("|"):
            cells = [c.strip() for c in s.strip("|").split("|")]
            if all(__import__("re").match(r"^:?-+:?$", c) for c in cells if c):
                in_table = True
                continue
            if cells:
                rows.append(cells)
    if rows and in_table:
        # 표 데이터 — 첫 행은 헤더로 사용
        headers = rows[0]
        data_rows = rows[1:11]
        n_cols = min(4, len(headers))
        n_rows = len(data_rows)
        if n_rows == 0:
            return
        col_w = (96 - 4) / n_cols
        row_h = max(4, 80 / (n_rows + 1) - 0.5)
        x0 = 4
        y_top = 88
        # 헤더
        for ci in range(n_cols):
            x = x0 + ci * col_w
            ax.add_patch(FancyBboxPatch(
                (x, y_top - row_h), col_w - 0.5, row_h,
                boxstyle="round,pad=0.2,rounding_size=0.3",
                facecolor=_hdr_face, edgecolor=_hdr_edge, linewidth=1.0,
            ))
            ax.text(x + col_w / 2, y_top - row_h / 2, str(headers[ci])[:18],
                    ha="center", va="center", fontsize=10, color="#fff", weight="bold")
        # 데이터
        for ri, row in enumerate(data_rows):
            y = y_top - (ri + 2) * row_h
            for ci in range(n_cols):
                x = x0 + ci * col_w
                cell = str(row[ci])[:24] if ci < len(row) else ""
                face = "#fff" if ri % 2 == 0 else "#f3f6fc"
                ax.add_patch(FancyBboxPatch(
                    (x, y), col_w - 0.5, row_h,
                    boxstyle="round,pad=0.15,rounding_size=0.2",
                    facecolor=face, edgecolor="#cccccc", linewidth=0.7,
                ))
                ax.text(x + 1, y + row_h / 2, cell,
                        ha="left", va="center", fontsize=9, color="#1e1e1e")
        return

    # 표가 아니면 카테고리 카드 그리드
    cat_lines = [ln.strip().lstrip("-*•").strip() for ln in lines_raw if ln.strip()]
    cat_lines = cat_lines[:8]
    if not cat_lines:
        return
    n = len(cat_lines)
    cols = 2 if n > 3 else 1
    rows_n = (n + cols - 1) // cols
    card_w = (96 - (cols - 1) * 3) / cols
    card_h = max(10, (80 - (rows_n - 1) * 2) / rows_n)
    x0, y_top = 4, 88
    card_colors = [
        ("#cfe2f3", "#3c78d8"),
        ("#d9ead3", "#6aa84f"),
        ("#fff2cc", "#bf9000"),
        ("#f4cccc", "#cc0000"),
        ("#ead1dc", "#a64d79"),
        ("#d9d2e9", "#674ea7"),
        ("#d0e0e3", "#45818e"),
        ("#e6e6e6", "#666666"),
    ]
    # palette가 있으면 카드 색을 palette 파생 (face,edge) 쌍으로 교체 (요구사항 7.3).
    if _fe:
        card_colors = _fe
    for i, txt in enumerate(cat_lines):
        r = i // cols
        c = i % cols
        x = x0 + c * (card_w + 3)
        y = y_top - (r + 1) * card_h - r * 2
        face, edge = card_colors[i % len(card_colors)]
        ax.add_patch(FancyBboxPatch(
            (x, y), card_w, card_h,
            boxstyle="round,pad=0.4,rounding_size=0.6",
            facecolor=face, edgecolor=edge, linewidth=1.4,
        ))
        if ":" in txt:
            head, comp = txt.split(":", 1)
            ax.text(x + card_w / 2, y + card_h - card_h * 0.3, head.strip()[:30],
                    ha="center", va="center", fontsize=11, color="#1e1e1e", weight="bold")
            comp_text = comp.strip()[:120]
            # 간단한 자동 줄바꿈
            words = comp_text.split(",")
            wrapped = []
            line = ""
            for w in words:
                w = w.strip()
                if len(line) + len(w) + 2 > 28:
                    if line:
                        wrapped.append(line)
                    line = w
                else:
                    line = (line + ", " + w) if line else w
            if line:
                wrapped.append(line)
            for li, lw in enumerate(wrapped[:3]):
                ax.text(x + card_w / 2, y + card_h - card_h * 0.55 - li * 2.2,
                        lw, ha="center", va="center", fontsize=9, color="#1e1e1e")
        else:
            ax.text(x + card_w / 2, y + card_h / 2, txt[:60],
                    ha="center", va="center", fontsize=10, color="#1e1e1e", weight="bold")


# === 섹션 의미 분류기 ===
# heading + body를 보고 어떤 다이어그램이 적합한지 결정. 같은 (kind, content) 조합은
# _force_generate_from_text에서 캐시되어 중복 생성을 방지한다.
def _mine_diagram_entities(text: str, max_items: int = 8) -> list:
    """프로즈/이미지 프롬프트에서 다이어그램 노드 라벨 후보를 추출한다.

    모델이 구조 슬라이드를 만들 때 본문(bullets) 없이 imagePrompt 문장만 주는 경우가
    많다(예: "folder tree of repository with frontend, backend, tests, docs, scripts").
    이때 콤마/슬래시/화살표/'and'/줄바꿈 등으로 분리해 짧은 명사형 토큰만 추려, 통짜
    래스터 대신 *편집 가능한* 네이티브 도형의 노드로 쓸 수 있게 한다.

    설명용 불용어(diagram/folder/blue/icon 등)는 제거하고, 4단어 초과·30자 초과·2자 미만
    토큰은 버린다. 중복 제거 후 최대 max_items개를 순서대로 반환한다.
    """
    if not text:
        return []
    import re as _re
    t = str(text)
    # 구분자 정규화 → 콤마
    t = t.replace("→", ",").replace("->", ",").replace("⇒", ",").replace("/", ",")
    t = _re.sub(r"\b(and|or|그리고|및|with|including|등|→|vs)\b", ",", t, flags=_re.IGNORECASE)
    for ch in ("\n", ";", "·", "•", "|", "·"):
        t = t.replace(ch, ",")
    STOP = {
        "diagram", "diagrams", "folder", "folders", "tree", "structure", "structures",
        "repository", "repo", "isometric", "blue", "icons", "icon", "showing", "show",
        "with", "of", "the", "a", "an", "style", "flat", "modern", "minimal", "clean",
        "vector", "illustration", "design", "layout", "view", "scene", "3d", "render",
        "구조", "폴더", "트리", "다이어그램", "아이콘", "구성", "구성도", "이미지", "image",
        "스타일", "벡터", "일러스트", "흐름도", "도식", "표현", "형태",
    }
    raw = [p.strip().strip(".-*:()[]{}\"'` \t").strip() for p in t.split(",")]
    out = []
    for p in raw:
        if not p:
            continue
        low = p.lower()
        if low in STOP:
            continue
        words = p.split()
        if len(words) > 4:
            continue
        if len(p) < 2 or len(p) > 30:
            continue
        # 순수 불용어 조합(예: "blue icons") 제거 — 모든 단어가 STOP이면 버림
        if all(w.lower() in STOP for w in words):
            continue
        if p not in out:
            out.append(p)
        if len(out) >= max_items:
            break
    return out


def _classify_section_diagram(heading: str, body: str, doc_title: str = "") -> tuple:
    """섹션의 적합한 다이어그램 종류와 콘텐츠를 결정.

    Returns:
        (kind, content) — kind는 "tree"/"flow"/"architecture"/"stack"/"block" 중 하나.
        kind가 빈 문자열이면 다이어그램 생성 안 함 (텍스트 전용 섹션).
    """
    text = f"{heading} {body}".lower()
    h = (heading or "").lower()

    # 1) 디렉토리/폴더/트리/구조 — 들여쓰기 라인이 많으면 tree
    #    heading 뿐 아니라 body/프롬프트(text)도 검사 — 모델이 본문 없이 imagePrompt만
    #    주는 구조 슬라이드(예: "folder tree of repository with frontend, backend ...")도 잡는다.
    _tree_kw = (
        "디렉토리", "디렉터리", "폴더", "tree", "구조", "directory", "folder",
        "structure", "repository", "repo", "hierarchy", "계층 구조", "조직도",
        "organization", "filesystem", "파일 구조", "디렉토리 구조", "폴더 구조",
    )
    if any(kw in text for kw in _tree_kw):
        # body에 들여쓰기 라인이 있는지
        indented = sum(1 for ln in (body or "").splitlines() if ln.startswith(("  ", "\t", "│", "├", "└")))
        if indented >= 2:
            return ("tree", body)
        # 들여쓰기 없으면 디렉토리 항목 추출 시도 (백틱 경로)
        items = []
        import re as _re
        for ln in (body or "").splitlines():
            s = ln.strip()
            if not s:
                continue
            for m in _re.finditer(r"`([^`]+)`", s):
                items.append(m.group(1))
        if items:
            tree_text = "\n".join(_path_to_tree_lines(items))
            return ("tree", tree_text)
        # 백틱 경로도 없으면 prose/프롬프트에서 엔티티(폴더명 등)를 마이닝해
        # 루트+자식 트리를 구성 → 통짜 래스터 대신 편집 가능한 도형으로.
        mined = _mine_diagram_entities(body)
        if len(mined) < 2:
            mined = _mine_diagram_entities(f"{heading} {body}")
        if len(mined) >= 2:
            root = (heading or doc_title or "Structure").strip()[:30] or "Structure"
            tree_text = root + "\n" + "\n".join("  " + m for m in mined)
            return ("tree", tree_text)
        # 마이닝 실패 — body가 있으면 그대로, 없으면 다이어그램 생성 포기(래스터 폴백)
        return ("tree", body) if (body or "").strip() else ("", "")

    # 1.4) 진행률/목표 달성(progress) / 지표·성과 대시보드(kpi) — 젠스파크 OKR 스타일.
    #      퍼센트 막대/KPI 카드는 일반 cards(#1.5)보다 메트릭 표현에 적합하므로 먼저 검사.
    _prog_kw = (
        "진척", "진행률", "달성", "달성률", "목표 대비", "okr", "마일스톤",
        "progress", "completion", "달성도", "이행률", "진행 현황",
    )
    _kpi_kw = (
        "지표", "성과", "kpi", "대시보드", "metric", "dashboard",
        "주요 수치", "핵심 지표", "실적", "수치",
    )
    _pct_count = (body or "").count("%")
    _has_num = any(c.isdigit() for c in (body or ""))
    # 진행률: 퍼센트 값이 2개 이상 + 진행/달성/OKR 키워드 → 막대 차트
    if any(kw in text for kw in _prog_kw) and _pct_count >= 2:
        return ("progress", body)
    # KPI: 지표/성과/대시보드 키워드 + 숫자 → 큰 숫자 카드. 단, _parse_kpis가
    #      값을 못 뽑으면 build_native_diagram이 False를 반환해 cards로 폴백됨.
    if any(kw in text for kw in _kpi_kw) and _has_num:
        return ("kpi", body)

    # 1.5) 지표/메트릭/특징 카드 — architecture 보다 *먼저* 검사(이슈4).
    #      heading에 메트릭/지표/성능 등 카드 신호가 있거나 ': ' 라인이 3개 이상이면
    #      cards로 분류. (이전: arch 키워드 '구성/시스템'이 먼저 매칭돼 메트릭
    #      슬라이드가 architecture→entity-mining으로 쉼표 분할되며 깨졌음.)
    _card_kw_h = ("지표", "메트릭", "metric", "kpi", "성과", "성능", "통계",
                  "수치", "특징", "장점", "핵심 기능", "항목", "feature",
                  "benefit", "highlight", "대시보드", "dashboard")
    _colon_lines0 = [ln for ln in (body or "").splitlines()
                     if ln.strip() and (": " in ln or "：" in ln)]
    if any(kw in h for kw in _card_kw_h) or len(_colon_lines0) >= 3:
        _cb0 = "\n".join(ln.strip().lstrip("-*•").strip()
                         for ln in (body or "").splitlines() if ln.strip())
        if len([ln for ln in _cb0.splitlines() if ln.strip()]) >= 2:
            return ("cards", _cb0)

    # 2) 시스템 아키텍처/계층/구성도 — architecture 다이어그램
    _arch_kw = (
        "아키텍처", "architecture", "계층", "tier", "system", "시스템",
        "구성도", "구성", "topology", "토폴로지", "infra", "인프라", "layered",
    )
    if any(kw in text for kw in _arch_kw):
        # body의 짧은 라인 중 레이어 키워드 우선
        layer_lines = []
        for ln in (body or "").splitlines():
            s = ln.strip().lstrip("-*•").strip()
            if not s:
                continue
            if any(kw in s.lower() for kw in (
                "계층", "tier", "layer", "프레젠테이션", "프론트", "백엔드", "데이터",
                "캐시", "presentation", "application", "data", "frontend", "backend",
            )):
                layer_lines.append(s)
        if not layer_lines:
            sentences = [s.strip() for s in (body or "").replace("\n", " ").split(".") if s.strip()]
            layer_lines = sentences[:5]
        # 라인이 빈약하면(프롬프트만 있는 경우) 엔티티 마이닝으로 보강
        if len([s for s in layer_lines if s]) < 2:
            mined = _mine_diagram_entities(body) or _mine_diagram_entities(f"{heading} {body}")
            if len(mined) >= 2:
                return ("architecture", "\n".join(mined[:6]))
            if not (body or "").strip():
                return ("", "")
        return ("architecture", "\n".join([s for s in layer_lines if s][:6]))

    # 3) 기술 스택 — stack 그리드
    if any(kw in h for kw in ("기술 스택", "tech stack", "stack", "스택", "기술")):
        # 마크다운 표가 있으면 그것을 우선
        return ("stack", body or "")

    # 4) 흐름/플로우/프로세스/순서 — flow
    if any(kw in text for kw in ("흐름", "플로우", "프로세스", "순서", "단계", "flow", "process", "step", "pipeline")):
        # body에서 -> 또는 → 패턴 추출
        flow_text = (body or "").replace("→", "->").replace("⇒", "->")
        if "->" in flow_text:
            return ("flow", flow_text)
        # 숫자 리스트나 불릿이 4개 이하면 flow
        bullets = [ln.strip().lstrip("-*•0123456789.").strip()
                   for ln in (body or "").splitlines() if ln.strip()]
        bullets = [b for b in bullets if b][:6]
        if 2 <= len(bullets) <= 6:
            return ("flow", " -> ".join(bullets))
        return ("block", body or "")

    # 5) 모듈/컴포넌트 분석 — block (수직 박스)
    if any(kw in h for kw in ("모듈", "module", "component", "컴포넌트", "분석")):
        bullets = []
        for ln in (body or "").splitlines():
            s = ln.strip().lstrip("-*•").strip()
            if s:
                # 첫 문장만 추출 (".으로 분리)
                first = s.split(".")[0].strip()
                if first:
                    bullets.append(first[:60])
        bullets = bullets[:6]
        if bullets:
            return ("block", "\n".join(bullets))
        return ("", "")

    # 5.5) 지표/메트릭/특징/항목 — '제목: 설명' 형태 카드 그리드(cards).
    #     KPI 불릿은 마침표가 없어 아래 fallback에서 누락되던 문제(이슈4) 해결.
    _card_kw = ("지표", "메트릭", "metric", "kpi", "성과", "성능", "통계",
                "수치", "특징", "장점", "핵심 기능", "항목", "feature",
                "benefit", "highlight")
    _colon_lines = [ln for ln in (body or "").splitlines()
                    if ln.strip() and (": " in ln or "：" in ln)]
    if any(kw in text for kw in _card_kw) or len(_colon_lines) >= 3:
        _cards_body = "\n".join(
            ln.strip().lstrip("-*•").strip()
            for ln in (body or "").splitlines() if ln.strip()
        )
        if len([ln for ln in _cards_body.splitlines() if ln.strip()]) >= 2:
            return ("cards", _cards_body)

    # 6) 그 외 — heading이 있고 body가 길면 block, 짧으면 다이어그램 없음
    if heading and len((body or "").strip()) > 100:
        # body의 핵심 문장 3-5개로 block
        sentences = [s.strip() for s in (body or "").replace("\n", " ").split(".") if s.strip()]
        if len(sentences) >= 2:
            return ("block", "\n".join(sentences[:5]))
    return ("", "")


# 진짜 "구조형 다이어그램"으로 보는 _classify_section_diagram kind 집합.
# 이 셋만 role="structural"(편집 가능 네이티브 도형 우선)으로 분류해 Req 3.1을 보존한다.
# kpi/cards/twocol/stack/block/progress 같은 "구조라기보다 고밀도 콘텐츠"는 content로 흡수해
# HTML 고밀도 레이아웃 경로로 보낸다.
_STRUCTURAL_DIAGRAM_KINDS = frozenset({"flow", "tree", "architecture"})


def _has_visual_intent(image_prompt: str, heading: str = "", body: str = "") -> bool:
    """슬라이드가 사진/일러스트(시각형) 이미지를 의도하는지 판정.

    visual intent = imagePrompt가 실제로 존재 **AND** 그 프롬프트가 구조 신호
    (화살표 체인 `A -> B`, 경로 토큰 `src/foo.js`, 마크다운 표)를 담고 있지 *않음*.
    구조 신호가 있는 프롬프트는 사진이 아니라 다이어그램 의도이므로 visual로 보지 않는다
    (`_looks_structural` 재사용). 이렇게 하면 "visual intent ∧ NOT kind → visual"
    규칙이 다이어그램형 프롬프트를 사진으로 오분류하지 않는다.
    """
    p = str(image_prompt or "").strip()
    if not p:
        return False
    try:
        if _looks_structural(p, heading, body):
            return False
    except Exception:
        # 구조 신호 판정 실패 시에는 보수적으로 imagePrompt 존재만으로 visual 인정.
        pass
    return True


def _hybrid_render_enabled(env: str) -> bool:
    """`AE_HYBRID_RENDER` opt-in 플래그를 파싱하는 순수 결정 함수.

    하이브리드 렌더 라우팅(pptx-ultra-quality-hybrid-render)의 활성화 게이트다.
    라우팅 로직은 포함하지 않으며 값 파싱만 담당한다.

    파싱 규칙 (기본 ON — A안: 하이브리드가 표준 동작, "0"은 킬스위치):
      - "0"                       → False (명시적 킬스위치 / legacy 렌더로 롤백)
      - "1"                       → True  (하이브리드 활성)
      - 미설정/""                 → True  (DEFAULT ON — 하이브리드가 표준 동작)
      - 그 외 인식 불가 값        → True + 경고 로그 1줄 (예: "2", "true", "on")

    불변식:
      - 결정론적: 동일 입력 → 동일 출력.
      - 어떤 입력(None/공백/유니코드/임의 문자열)에도 raise하지 않는다.
      - LLM/네트워크/게이트웨이 호출 없음.

    _Requirements: 1.7, 6.1, 6.2, 6.5 / Design: Property 20_
    """
    # None 이나 문자열이 아닌 입력도 방어적으로 문자열화 (raise 금지).
    try:
        raw = "" if env is None else str(env)
    except Exception:
        return True
    val = raw.strip()
    if val == "0":
        return False
    if val in ("", "1"):
        return True
    # 인식 불가 값 → 기본 ON 처리 + 경고 1줄(≤200자). 킬스위치는 "0"만.
    print(f"[HybridRender] AE_HYBRID_RENDER 인식 불가 값 '{val[:64]}' — 기본 ON으로 처리(킬스위치는 '0')")
    return True


def _classify_slide_role(slide: dict, is_cover: bool, doc_title: str = "",
                         *, bg_has_baked_text: bool = False) -> str:
    """슬라이드의 역할(role)을 결정한다 — 주 렌더러와 Vertex 이미지 사용처의 핵심 입력.

    반환값: ``cover | section | structural | content | visual`` 중 하나(문자열).

    결정 규칙 (design Fix Implementation §1 의 classifyRole):
      - ``is_cover`` → ``cover``
      - ``_classify_section_diagram`` kind ∈ {flow, tree, architecture} → ``structural``
        (진짜 구조형만 structural — Req 3.1 보존)
      - visual intent(imagePrompt 존재 ∧ 구조 신호 아님) ∧ NOT kind → ``visual``
      - 그 외(kpi/cards/twocol 등 고밀도 콘텐츠) → ``content``

    LLM/게이트웨이 추가 호출 없음 — 기존 결정론적 분류기(`_classify_section_diagram`,
    `_looks_structural`)만 재사용한다(Bedrock Gateway-only 제약 무관).

    note: ``section`` 은 role 열거형의 일원이지만 현 휴리스틱 규칙에서는 산출하지 않는다
    (주요 섹션 구분 슬라이드 등 향후 확장을 위해 예약). 반환 타입에는 포함된다.
    """
    if is_cover:
        return "cover"

    if not isinstance(slide, dict):
        slide = {}

    heading = str(slide.get("title") or slide.get("heading") or "")
    bullets = slide.get("bullets") or []
    if not isinstance(bullets, list):
        bullets = []
    _body_lines = [str(b) for b in bullets if str(b).strip()]
    body = "\n".join(_body_lines)
    _extra_body = slide.get("body")
    if isinstance(_extra_body, str) and _extra_body.strip():
        body = (body + "\n" + _extra_body).strip() if body else _extra_body.strip()
    img_prompt = str(slide.get("imagePrompt") or "").strip()

    # 1) 다이어그램 kind 분류 — 슬라이드의 실제 콘텐츠(title/body)로만 판정한다.
    #    (design Fix Implementation §1 classifyRole: kind := _classify_section_diagram(title, body)).
    #    imagePrompt 구조 키워드를 kind 로 흡수하면 비구조형 content 슬라이드가 structural 로
    #    오분류되어 Vertex 사전생성이 억제(손실)되므로 폴백하지 않는다.
    try:
        kind, _ = _classify_section_diagram(heading, body, doc_title)
    except Exception:
        kind = ""

    # 2) 진짜 구조형 kind 만 structural — 편집 가능 네이티브 도형 우선(Req 3.1 보존)
    if kind in _STRUCTURAL_DIAGRAM_KINDS:
        return "structural"

    # 3) 시각형(사진/일러스트) — imagePrompt가 있고 구조 신호가 아니며 diagram kind도 아님
    if not kind and _has_visual_intent(img_prompt, heading, body):
        # 결함 B (task 3.4 / design 3a): 풀블리드 배경에 텍스트가 구워져 있으면 그
        # 이미지를 슬라이드 본문 비주얼 캐리어(visual)로 쓰지 않는다 — content 로 강등해
        # 편집 가능 콘텐츠 경로로 보낸다(생성 이미지는 _select_render_plan 에서 backdrop
        # 슬롯으로 보존되어 손실-0 유지). bg_has_baked_text=False(기본)이면 기존과
        # 동일하게 visual 반환(바이트 보존).
        if bg_has_baked_text:
            return "content"
        return "visual"

    # 4) 그 외 — 고밀도 콘텐츠(HTML 고밀도 레이아웃 대상)
    return "content"


def _select_render_plan(*, has_vertex_image: bool, has_native_diagram: bool,
                        has_image_file: bool, has_slide_bg: bool,
                        role: str, html_enabled: bool,
                        bg_has_baked_text: bool = False) -> dict:
    """슬라이드 미디어 상태 → 손실 없는 렌더 플랜(design Fix Implementation §3).

    폐기형 임베드 가드(`if not native_diag and not img_file and not slide_bg:
    img_file = pre`)를 대체하는 결정 규칙. 슬라이드마다 정확히 하나의 *주 렌더러*와
    생성된 Vertex 이미지(`pre`)의 *사용처(slot)*를 정한다.

    반환::

        {"primary": "HTML" | "NATIVE_SHAPES" | "VERTEX_IMAGE",
         "vertex_slot": "hero" | "backdrop" | "visual" | "none",
         "body_separated": bool}

    분기 규칙 (design Fix Implementation §3):
      - slide_bg 존재          → primary=HTML,          slot=hero
        (HTML 풀블리드가 주 렌더러. pre는 히어로/이미지 슬롯으로 합성하거나
         합성 불가 시 on-slide 레이어로 보존.)
      - native_diagram 존재    → primary=NATIVE_SHAPES, slot=backdrop
        (네이티브 도형 유지. pre는 _native_over_bg/_eff_bg backdrop으로 보존.)
      - image_file 존재        → primary=VERTEX_IMAGE,  slot=backdrop
        (caller가 명시한 imageFile이 슬라이드 비주얼 — 우선순위 보존. pre는 backdrop.)
      - 그 외(cover/visual/비구조 content) → primary=VERTEX_IMAGE, slot=visual
        (이미지가 곧 슬라이드 비주얼: img_file := pre. 기존 가드 동작과 동일.)
      - content ∧ HTML 비활성   → 네이티브 카드 경로(native_diagram이 부여되면 위
        분기에서 흡수)이며 pre는 backdrop으로 보존된다.

    핵심 불변식(손실 0): ``has_vertex_image`` 이면 반환 slot 은 결코 ``"none"`` 이 아니다 —
    생성된 Vertex 이미지는 어떤 분기에서도 폐기되지 않는다. 명시 우선순위 보존: caller가
    지정한 imageFile/slideBackground 경로가 항상 주 렌더러로 유지된다.

    이 함수는 LLM/게이트웨이 호출이 없는 순수 결정 함수다(전역 정의·결정성, design Property 3).
    """
    if has_slide_bg:
        primary = "HTML"
        slot = "hero"
    elif has_native_diagram:
        primary = "NATIVE_SHAPES"
        slot = "backdrop"
    elif has_image_file:
        primary = "VERTEX_IMAGE"
        slot = "backdrop"
    else:
        # cover/visual/비구조 content — 이미지가 곧 슬라이드 비주얼(기존 가드 동작 보존).
        primary = "VERTEX_IMAGE"
        slot = "visual"

    # --- 결함 B 결정 규칙 보강 (spec: pptx-overlay-collision-fix, task 3.4 / design 3a) ---
    # bg_has_baked_text=False(기본, 기존 모든 호출)이면 이 블록은 동작하지 않아 위에서 정한
    # primary/vertex_slot 이 그대로 유지된다(바이트 보존). 참일 때만 본문 캐리어 차단.
    body_separated = False
    if bg_has_baked_text:
        # 생성 이미지가 본문 캐리어(슬라이드 비주얼 visual / HTML 풀블리드 본문 hero)가
        # 되는 상황인지.
        would_carry_body = slot in ("visual", "hero")
        if role == "structural":
            # 구조형: 편집 가능 네이티브 도형을 주 렌더러로 유지, 구워진-텍스트 이미지는
            # 본문 캐리어로 쓰지 않고 backdrop 으로만 보존(손실-0).
            primary = "NATIVE_SHAPES"
            if has_vertex_image:
                slot = "backdrop"
            body_separated = True
        elif would_carry_body:
            # 비구조라도 구워진-텍스트 풀블리드가 본문 캐리어가 되는 상황 → 본문/배경 분리
            # 의도 신호(좌표 분리는 task 3.5의 body_safe_area 합성 경로에서 적용).
            body_separated = True

    # 손실 0 불변식 강제: Vertex 이미지가 없으면 slot은 의미 없음("none").
    if not has_vertex_image:
        slot = "none"
    return {"primary": primary, "vertex_slot": slot, "body_separated": body_separated}


def _select_hybrid_render_plan(*, role: str, vertex_enabled: bool,
                               html_enabled: bool, has_vertex_image: bool,
                               has_native_diagram: bool, has_image_file: bool,
                               has_slide_bg: bool) -> dict:
    """슬라이드 역할 × Vertex/HTML 상태 → 하이브리드 렌더 플랜(design 결정 테이블).

    `_tool_generate_pptx`의 하이브리드 라우팅 레이어가 **caller 미지정 슬라이드**에 대해
    역할 기반 *주 렌더러*를 배정하기 위해 호출하는 순수 결정 함수다. caller가
    ``imageFile``/``slideBackground``를 명시한 슬라이드는 이 함수를 거치지 않고 항상
    기존 :func:`_select_render_plan` 경로에 위임된다(caller 우선순위 보존, R4.2).

    반환::

        {"primary": "VERTEX_FULLBLEED" | "HTML_EDITABLE" | "NATIVE_EDITABLE" | "NATIVE_SHAPES",
         "vertex_slot": "visual" | "hero" | "backdrop" | "none",
         "editable": bool}

    결정 테이블 (design "결정 테이블" 절 그대로):
      - ``cover``/``section``/``visual`` ∧ vertex_enabled
            → primary=VERTEX_FULLBLEED, slot=visual, editable=False  (R1.2)
      - ``cover``/``section``/``visual`` ∧ ¬vertex_enabled
            → primary=HTML_EDITABLE(html on) / NATIVE_EDITABLE(html off),
              slot=none, editable=True                                (R1.6 편집 폴백)
      - ``content`` (임의 vertex/html)
            → primary=NATIVE_EDITABLE, slot=hero if has_vertex_image else none,
              editable=True                                           (R1.3/R2.1/R2.2)
      - ``structural`` (임의 vertex/html)
            → primary=NATIVE_SHAPES, slot=backdrop if has_vertex_image else none,
              editable=True                                           (R1.4/R4.5)

    R1.8(방어적 이중화): role이 ``{cover, section, structural, content, visual}`` 밖의
    모호/미정의/복수 후보 값이면 진입 시점에 결정론적으로 ``content``로 확정한다.
    (:func:`_classify_slide_role`도 이미 ``content``로 폴백하므로 이중 안전망.)

    불변식 (design Data Models §RenderPlan):
      - ``primary``는 정확히 1개(주 렌더러 유일성, R1.5).
      - ``has_vertex_image == True ⇒ vertex_slot != "none"``(손실-0, R4.1). 최종 게이트는
        기존 :func:`_select_render_plan`이 담당하지만, 여기서도 slot이 ``"none"``이면
        ``"backdrop"``으로 승격해 생성 이미지를 폐기하지 않는다.
      - ``primary == "VERTEX_FULLBLEED" ⇒ role ∈ {cover, section, visual} ∧ vertex_enabled``(R1.2).
      - ``role == "content" ⇒ editable ∧ primary == "NATIVE_EDITABLE"``(R2.1/R2.2).
      - ``role == "structural" ⇒ primary == "NATIVE_SHAPES"``(R1.4/R4.5).

    ``has_native_diagram``/``has_image_file``/``has_slide_bg``는 시그니처 완전성 및
    :func:`_select_render_plan` 위임 검증을 위해 수용한다. caller 지정 미디어
    (image_file/slide_bg) 우선순위는 호출부(와이어링)에서 처리되므로 결정 테이블 자체는
    역할 기반이다.

    LLM/게이트웨이/네트워크 호출이 없는 순수 결정 함수다(결정성 보장, R3.7/R4.4).
    """
    # R1.8 — 모호/미정의/복수 후보 role은 결정론적으로 content 확정(방어적 이중화).
    _VALID_ROLES = ("cover", "section", "structural", "content", "visual")
    if role not in _VALID_ROLES:
        role = "content"

    _FULLBLEED_ROLES = ("cover", "section", "visual")

    if role in _FULLBLEED_ROLES:
        if vertex_enabled:
            # R1.2 — Vertex 초고품질 풀블리드가 주 렌더러(편집 불가).
            primary = "VERTEX_FULLBLEED"
            slot = "visual"
            editable = False
        else:
            # R1.6 — Vertex 비활성 → 편집 가능 고밀도로 강등 폴백.
            primary = "HTML_EDITABLE" if html_enabled else "NATIVE_EDITABLE"
            slot = "none"
            editable = True
    elif role == "structural":
        # R1.4/R4.5 — 편집 가능 네이티브 도형 유지, 생성 이미지는 backdrop 보존(래스터화 0).
        primary = "NATIVE_SHAPES"
        slot = "backdrop" if has_vertex_image else "none"
        editable = True
    else:
        # role == "content" — R1.3/R2.1/R2.2 — 항상 편집 가능 네이티브 고밀도.
        # Vertex 히어로가 있으면 바운디드 Image_Slot(hero)에 합성, 없으면 네이티브만.
        primary = "NATIVE_EDITABLE"
        slot = "hero" if has_vertex_image else "none"
        editable = True

    # 손실-0 안전망(R4.1) — 생성된 Vertex 이미지는 어떤 경로에서도 폐기하지 않는다.
    # 최종 검증은 _select_render_plan에 위임하되, 여기서도 slot=="none"으로 남으면
    # backdrop으로 승격해 has_vertex_image ⇒ slot != "none" 불변식을 유지한다.
    if has_vertex_image and slot == "none":
        slot = "backdrop"

    return {"primary": primary, "vertex_slot": slot, "editable": editable}


# content 경로에서 히어로 이미지를 "슬롯"으로 합성할 수 있는 네이티브 레이아웃 집합.
# 이 레이아웃들은 우측 컬럼(이미지 컬럼) 성격의 바운디드 영역을 가진다(설계: two_column
# 우측 / objective_detail 이미지 컬럼). 그 외 레이아웃(section_divider/timeline/
# architecture/cover/fallback 등)은 슬롯을 호스팅하지 못하므로 바운디드 on-slide 레이어로
# 보존한다(R2.4, 손실-0). 어느 경우든 슬라이드 전체를 덮는 풀블리드 PICTURE는 만들지 않는다.
_CONTENT_SLOT_LAYOUTS: frozenset = frozenset({"two_column", "feature_grid", "comparison"})

# content 히어로 합성용 바운디드 영역(인치) — 모두 슬라이드(13.333x7.5)보다 작아
# 절대 풀블리드 PICTURE가 되지 않는다. SLOT=우측 이미지 컬럼, PRES=바운디드 보존 영역.
_CONTENT_SLOT_REGION = (8.45, 1.85, 4.25, 4.20)   # 우측 이미지 컬럼(바운디드)
_CONTENT_PRES_REGION = (8.90, 4.05, 3.85, 2.95)   # 바운디드 보존(우하단, back-most)


def _render_content_editable(slide, prs, data, tokens, hero_rel, palette):
    """content 슬라이드를 편집 가능 네이티브 경로로 조립 렌더한다 (task 5.2).

    HTML이 전역 on이어도 :func:`_generate_html_slide_for_section`의 풀블리드 PNG 바이크를
    **우회**하고 ``native_layout_renderer.render_native_layout``(편집 가능 네이티브 도형/
    텍스트)로 렌더한다(R2.2). Vertex 히어로/액센트 이미지(``hero_rel``)가 있으면 슬라이드
    전체보다 작은 **바운디드 Image_Slot**에 ``add_picture(rel, left, top, width, height)``로
    합성한다(R2.1/R2.3). 레이아웃이 바운디드 슬롯을 호스팅할 수 없거나 합성이 실패해도
    이미지를 **바운디드 on-slide 레이어**(back-most, 여전히 비풀블리드)로 보존해 폐기하지
    않는다(R2.4, 손실-0). ``add_picture`` 실패는 예외를 전파하지 않고 보존 폴백으로 넘어간다.

    측정 가능한 "편집 가능" 정의(설계)::

        편집 가능 텍스트 run 개수 >= 1  AND  슬라이드 전체(13.333in x 7.5in) PICTURE 개수 == 0

    ``native_layout_renderer.maybe_add_decorative_background``(풀블리드 장식 배경)는
    content 경로에서 **사용하지 않는다** — 풀블리드 PICTURE를 만들어 R2.1을 깨기 때문이다.
    (render_native_layout는 옵트인 OFF/자격증명 부재 기본값으로 호출하므로 장식 배경이
    방출되지 않는다.)

    Args:
        slide: python-pptx Slide (이미 add_slide로 생성됨)
        prs: python-pptx Presentation
        data: render_native_layout/slide_templates 데이터 스키마(dict). 선택적 ``layout``
            키로 레이아웃 힌트를 줄 수 있고, ``bullets``로 고밀도 후보를 유도한다.
        tokens: design_tokens_for_profile 결과(색/여백/타이포)
        hero_rel: Vertex 히어로/액센트 이미지 로컬 경로(있으면 바운디드 합성/보존)
        palette: 선택 팔레트(``[#RRGGBB, ...]`` 또는 None)

    Returns:
        dict::

            {"ok": bool,               # 편집 가능 네이티브 방출 성공 여부
             "layout": str,            # 사용된 레이아웃명(또는 "fallback")
             "editable": True,         # content 경로는 항상 편집 가능 의도
             "title_count": int,       # 방출된 제목 셰이프 수(0 또는 1)
             "image_placed": bool,     # 히어로를 바운디드 슬롯에 합성했는지
             "image_preserved": bool}  # 히어로를 바운디드 보존 레이어로 보존했는지

        LLM/네트워크 호출 없음(순수 네이티브 방출 + 로컬 이미지 임베드).

    NOTE: 본 함수는 조립 프리미티브다. ``_tool_generate_pptx`` 배선은 task 7.4에서 수행한다.
    """
    # 렌더러 import (실패해도 예외 전파 금지 — 손실-0 보존).
    try:
        from ai_engine.native_layout_renderer import (
            map_to_native_layout as _m2nl,
            render_native_layout as _rnl,
            render_native_fallback as _rnf,
        )
    except ImportError:
        from native_layout_renderer import (
            map_to_native_layout as _m2nl,
            render_native_layout as _rnl,
            render_native_fallback as _rnf,
        )

    _data = data if isinstance(data, dict) else {}
    _tokens = tokens if isinstance(tokens, dict) else {}

    # 제목/불릿 추출(방어적) — 합성 레이아웃 후보 구성에 사용.
    try:
        _t = _data.get("title")
        _title = str(_t).strip() if _t is not None else ""
    except Exception:
        _title = ""
    _bullets = []
    try:
        _b = _data.get("bullets")
        if isinstance(_b, (list, tuple)):
            _bullets = [str(x).strip() for x in _b if str(x).strip()]
    except Exception:
        _bullets = []

    # --- 1) 편집 가능 네이티브 렌더 (풀블리드 PNG 바이크 우회, R2.2) ---
    #     레이아웃 후보 사다리: data의 layout 힌트 -> 합성(feature_grid/two_column) -> 최소.
    #     각 후보를 render_native_layout 로 시도(ok=True 면 채택). 어떤 경우에도 통짜
    #     이미지로 굽지 않는다(편집 가능 run >= 1 보장은 마지막 fallback 이 담당).
    _cands = []
    _hint = ""
    try:
        if _data.get("layout"):
            _hint = _m2nl(str(_data.get("layout") or ""))
    except Exception:
        _hint = ""
    if _hint:
        _cands.append((_hint, _data))
    if _bullets:
        _cands.append(("feature_grid",
                       {"title": _title,
                        "features": [{"title": _bx, "description": ""} for _bx in _bullets]}))
        _half = (len(_bullets) + 1) // 2
        _cands.append(("two_column",
                       {"title": _title,
                        "left_content": _bullets[:_half] or [_title or "*"],
                        "right_content": _bullets[_half:] or [_title or "*"]}))
    _cands.append(("section_divider", {"title": _title or "Slide"}))

    _used_layout = ""
    _title_count = 0
    _native_ok = False
    for _cl, _cd in _cands:
        try:
            _res = _rnl(slide, prs, _cl, _cd, _tokens, palette=palette)
        except Exception:
            _res = None
        if _res is not None and getattr(_res, "ok", False):
            _native_ok = True
            _used_layout = _cl
            _title_count = int(getattr(_res, "title_count", 0) or 0)
            break

    # 리치 후보 전부 실패 -> 편집 가능 텍스트 폴백(통짜 금지, 편집 run >= 1 보장, 손실-0).
    if not _native_ok:
        try:
            _fb = _rnf(slide, _data, _tokens)
            _native_ok = bool(getattr(_fb, "ok", False))
            _used_layout = "fallback"
            _title_count = int(getattr(_fb, "title_count", 0) or 0)
        except Exception:
            _native_ok = False

    # --- 2/3) Vertex 히어로 바운디드 합성/보존 (풀블리드 PICTURE 절대 금지, 손실-0) ---
    def _try_bounded_picture(_region, *, back_most):
        """히어로를 바운디드 영역에 비율 보존 fit 하여 add_picture. 성공 시 True.

        영역/최종 rect는 슬라이드(13.333x7.5)보다 작게 유지되므로 풀블리드 PICTURE가
        되지 않는다. 실패 시 예외를 전파하지 않고 False 를 반환한다(보존 폴백 대상).
        """
        try:
            from pptx.util import Inches as _In
        except Exception:
            return False
        region_l, region_t, region_w, region_h = _region
        iw, ih = 4, 3
        try:
            from PIL import Image as _PILc
            with _PILc.open(hero_rel) as _im:
                iw, ih = _im.size
        except Exception:
            pass
        ar = (iw / ih) if ih else 1.3333
        draw_w = region_w
        draw_h = (region_w / ar) if ar else region_h
        if draw_h > region_h:
            draw_h = region_h
            draw_w = region_h * ar
        # 영역 내 중앙정렬(영역보다 커지지 않으므로 항상 바운디드).
        off_l = region_l + (region_w - draw_w) / 2.0
        off_t = region_t + (region_h - draw_h) / 2.0
        # 경계 클램프(있으면) — 영역이 이미 슬라이드 안이라 통상 no-op, 안전망.
        try:
            try:
                from layout_geometry import clamp_into_bounds as _clampb
            except Exception:
                from ai_engine.layout_geometry import clamp_into_bounds as _clampb
            off_l, off_t, draw_w, draw_h = _clampb((off_l, off_t, draw_w, draw_h))
        except Exception:
            pass
        try:
            _pic = slide.shapes.add_picture(
                hero_rel, _In(off_l), _In(off_t),
                width=_In(draw_w), height=_In(draw_h),
            )
        except Exception:
            return False
        if back_most:
            # 콘텐츠 텍스트보다 뒤(back-most)로 이동 — 여전히 비풀블리드.
            try:
                _spTree = _pic._element.getparent()
                _spTree.remove(_pic._element)
                _spTree.insert(2, _pic._element)  # 0=nvGrpSpPr, 1=grpSpPr, 2+=shapes
            except Exception:
                pass
        return True

    _image_placed = False
    _image_preserved = False
    if hero_rel:
        if _used_layout in _CONTENT_SLOT_LAYOUTS:
            # 슬롯 호스팅 레이아웃 -> 우측 이미지 컬럼(바운디드)에 합성.
            _image_placed = _try_bounded_picture(_CONTENT_SLOT_REGION, back_most=False)
            if not _image_placed:
                # add_picture 실패 -> 예외 전파 없이 바운디드 보존 폴백(R2.4).
                _image_preserved = _try_bounded_picture(_CONTENT_PRES_REGION, back_most=True)
        else:
            # 슬롯 미지원 레이아웃 -> 바운디드 on-slide 레이어로 보존(폐기 금지, R2.4).
            _image_preserved = _try_bounded_picture(_CONTENT_PRES_REGION, back_most=True)

    return {
        "ok": _native_ok,
        "layout": _used_layout,
        "editable": True,
        "title_count": _title_count,
        "image_placed": _image_placed,
        "image_preserved": _image_preserved,
    }


def _path_to_tree_lines(paths: list) -> list:
    """`/src/components`, `/src/services` 같은 경로 리스트를 들여쓰기 트리로 변환."""
    if not paths:
        return []
    # 정규화 + 정렬
    norm = []
    for p in paths:
        p = p.strip()
        if not p:
            continue
        p = p.lstrip("/").rstrip("/")
        if p:
            norm.append(p)
    norm = sorted(set(norm))
    out = []
    seen = set()
    for p in norm:
        parts = p.split("/")
        for d in range(len(parts)):
            sub = "/".join(parts[: d + 1])
            if sub in seen:
                continue
            seen.add(sub)
            indent = "  " * d
            name = parts[d]
            # 마지막 segment가 뒤에 / 안 붙은 plain이면 디렉토리로 가정 (깊이 있는 항목)
            is_dir = d < len(parts) - 1 or "." not in name
            out.append(f"{indent}{name}{'/' if is_dir else ''}")
    return out


def _strip_emoji(text, keep_if_empty=True):
    """문자열에서 일반 이모지/픽토그램/딩뱃 글리프를 제거한다.

    사용자 정책: 슬라이드/문서 텍스트에 일반 이모지 아이콘(📁📊✅🚀 등)을 절대 쓰지 않는다.
    한글·영문·숫자·일반 문장부호와 흐름 화살표(→←↑↓)는 보존하고 이모지 블록만 제거한다.
    전부 제거되어 빈 문자열이 되면 원본을 보존한다.
    """
    if not isinstance(text, str) or not text:
        return text
    import re as _re
    _emoji = _re.compile(
        "["
        "\U0001F300-\U0001FAFF"
        "\U00002600-\U000027BF"
        "\U00002B00-\U00002BFF"
        "\U0001F1E6-\U0001F1FF"
        "\U0000FE00-\U0000FE0F"
        "\U0000200D"
        "\U00002190-\U000021FF"
        "]",
        flags=_re.UNICODE,
    )
    _keep = {"\u2192", "\u2190", "\u2191", "\u2193"}
    cleaned = "".join("" if (_emoji.match(ch) and ch not in _keep) else ch for ch in text)
    cleaned = _re.sub(r"[ \t]{2,}", " ", cleaned).strip()
    return cleaned if (cleaned or not keep_if_empty) else text


def _normalize_doc_input(tool_input: dict, default_kind: str = "sections"):
    """모델이 도구 input의 shape을 헷갈려서 보내도 정상 형태로 정규화.

    PDF는 sections, PPTX는 slides, XLSX는 sheets, DOCX는 sections를 기대하지만
    LLM이 종종 다른 키 이름이나 잘못된 구조로 보냄. 이 함수가 유연하게 받아
    공통 형태로 변환한다.

    Args:
        tool_input: 모델이 보낸 raw input dict
        default_kind: 'sections' | 'slides' | 'sheets'

    Returns:
        (title: str, items: list[dict])
        items의 형태:
        - sections: [{heading, body, bullets, level?}]
        - slides:   [{title, bullets, body?}]
        - sheets:   [{name, headers, rows}]
    """
    if not isinstance(tool_input, dict):
        return "", []

    title = (
        tool_input.get("title")
        or tool_input.get("name")
        or tool_input.get("heading")
        or tool_input.get("file_title")
        or "Untitled"
    )
    title = _strip_emoji(str(title).strip())

    # 후보 키 — LLM이 다양한 이름으로 보낼 수 있음
    candidate_keys = ["sections", "slides", "sheets", "pages", "items",
                      "content", "data", "body", "rows", default_kind]
    raw_items = None
    for k in candidate_keys:
        v = tool_input.get(k)
        if isinstance(v, list) and v:
            raw_items = v
            break

    # 단일 dict로 보낸 경우 (예: {"section": {...}})
    if raw_items is None:
        for k in ("section", "slide", "sheet", "page"):
            v = tool_input.get(k)
            if isinstance(v, dict):
                raw_items = [v]
                break

    # 문자열로 본문만 보낸 경우 — 단일 섹션으로 wrap
    if raw_items is None:
        body = tool_input.get("body") or tool_input.get("text") or tool_input.get("description")
        if body:
            return title, [{"heading": title, "body": str(body)}]
        return title, []

    # 각 item 정규화
    out = []
    for item in raw_items:
        if not isinstance(item, dict):
            # 문자열 item — 본문으로 처리
            out.append({"heading": "", "body": str(item)})
            continue
        n = dict(item)  # shallow copy
        # Common: heading aliases
        if "heading" not in n:
            n["heading"] = n.get("title") or n.get("name") or n.get("header") or ""
        # Common: body aliases
        if "body" not in n:
            n["body"] = n.get("content") or n.get("text") or n.get("description") or ""
        # Common: bullets aliases
        if "bullets" not in n:
            n["bullets"] = n.get("points") or n.get("items") or n.get("list") or []
        if not isinstance(n.get("bullets"), list):
            n["bullets"] = []
        # Common: imagePrompt aliases (canonical = imagePrompt)
        # PDF/DOCX/PPTX 모두 동일한 alias 집합을 인식 — LLM이 다양한 키 이름으로 보내도 보존
        if not n.get("imagePrompt"):
            for alias in ("image_prompt", "imageDescription", "image_description",
                          "image", "picture"):
                v = n.get(alias)
                if isinstance(v, str) and v.strip():
                    n["imagePrompt"] = v.strip()
                    break
        # 비-문자열/빈 문자열은 제거하여 downstream에서 단순화
        if not isinstance(n.get("imagePrompt"), str) or not n.get("imagePrompt", "").strip():
            n.pop("imagePrompt", None)
        # imageFile (사전 렌더된 이미지 경로) — Bedrock 호출 없이 바로 임베드
        # 호출자(_force_generate_from_text)가 native matplotlib로 미리 그려놓은 PNG를 패스할 때 사용
        if not n.get("imageFile"):
            for alias in ("image_file", "imagePath", "image_path", "img_file", "img_path"):
                v = n.get(alias)
                if isinstance(v, str) and v.strip():
                    n["imageFile"] = v.strip()
                    break
        if not isinstance(n.get("imageFile"), str) or not n.get("imageFile", "").strip():
            n.pop("imageFile", None)
        # slideBackground (HTML→PNG full-bleed slide capture, Genspark/Gamma class).
        # 존재 시 PPTX는 슬라이드 전체(left=0, top=0, 13.33×7.5 in)를 이 이미지로 덮음.
        # PDF/DOCX는 imageFile과 동일하게 max-cap 적용 (페이지가 가려지지 않도록).
        if not n.get("slideBackground"):
            for alias in ("slide_background", "background", "bg", "fullBleed", "full_bleed"):
                v = n.get(alias)
                if isinstance(v, str) and v.strip():
                    n["slideBackground"] = v.strip()
                    break
        if not isinstance(n.get("slideBackground"), str) or not n.get("slideBackground", "").strip():
            n.pop("slideBackground", None)
        # Slides: ensure title (PPTX uses 'title' as slide title)
        if default_kind == "slides":
            n["title"] = n.get("title") or n.get("heading") or ""
        # Sheets: ensure name/headers/rows
        if default_kind == "sheets":
            n["name"] = n.get("name") or n.get("title") or n.get("sheet") or "Sheet1"
            n["headers"] = n.get("headers") or n.get("columns") or n.get("header") or []
            if not isinstance(n["headers"], list):
                n["headers"] = []
            n["rows"] = n.get("rows") or n.get("data") or []
            if not isinstance(n["rows"], list):
                n["rows"] = []
        out.append(n)

    # 일반 이모지 제거 — 제목/헤딩/본문/불릿에서 이모지 글리프 제거(이미지 프롬프트·시트 데이터 제외).
    for _n in out:
        for _k in ("heading", "title", "body"):
            if isinstance(_n.get(_k), str):
                _n[_k] = _strip_emoji(_n[_k])
        if isinstance(_n.get("bullets"), list):
            _bs = [_strip_emoji(str(_b), keep_if_empty=False) for _b in _n["bullets"]]
            _n["bullets"] = [_b for _b in _bs if _b and _b.strip()]

    return title, out


async def _tool_generate_pdf(tool_input: dict, project_path: str, aws_profile: str = '', bedrock_user: str = '') -> str:  # [patched-credentials]
    """Generate a PDF document using reportlab. Accepts lenient input shapes.

    각 섹션은 선택적으로 imagePrompt를 가질 수 있으며, 존재 시
    _tool_generate_image로 이미지를 생성하고 heading 다음, body 이전에
    Image flowable로 임베드한다 (실패 시 텍스트만으로 진행).
    """
    title, sections = _normalize_doc_input(tool_input, default_kind="sections")

    if not title:
        return json.dumps({"error": "title is required"})
    if not sections:
        return json.dumps({"error": "sections is required"})

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image
        from reportlab.lib.units import cm
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "reportlab", "hint": "pip install reportlab"})

    import time as _t, re as _re
    # Always use a locally-existing directory for generated media.
    # Remote project_path may not exist locally; fall back to cwd.
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _slug_from_title(title) or "doc"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.pdf"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    try:
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()

        # 한글 폰트 등록 — reportlab 기본 Helvetica는 한글 글리프가 없어
        # PDF에 한글이 ■■■ 박스로 깨진다. 시스템 폰트(AppleSDGothic/Noto/Malgun)를
        # 자동 탐지하여 모든 스타일에 적용. 폰트가 없으면 Helvetica fallback.
        _kfont = _register_korean_font_for_reportlab()
        # Title은 기본 24pt — 한글 폰트로 교체
        styles["Title"].fontName = _kfont
        # Req 4.2: enforce Heading2 = 14pt bold, Normal = 10pt explicitly so
        # the contract holds regardless of reportlab version defaults.
        # 한글 폰트는 보통 단일 weight라 bold도 같은 폰트 사용 (registerFontFamily로 매핑됨).
        styles["Heading2"].fontName = _kfont
        styles["Heading2"].fontSize = 14
        styles["Normal"].fontName = _kfont
        styles["Normal"].fontSize = 10
        # Heading1/3도 한글 폰트로 (있을 때만)
        for _sn in ("Heading1", "Heading3", "BodyText"):
            try:
                styles[_sn].fontName = _kfont
            except KeyError:
                pass
        story = [Paragraph(title, styles["Title"]), Spacer(1, 1 * cm)]

        # 이미지 폭/높이 캡 — 한 장이 전체 페이지를 차지하지 않도록
        max_img_w = float(doc.width)
        max_img_h = 12 * cm

        for sec in sections:
            heading = sec.get("heading", "") if isinstance(sec, dict) else ""
            body = sec.get("body", "") if isinstance(sec, dict) else str(sec)
            img_prompt = sec.get("imagePrompt", "") if isinstance(sec, dict) else ""
            img_file = sec.get("imageFile", "") if isinstance(sec, dict) else ""
            slide_bg = sec.get("slideBackground", "") if isinstance(sec, dict) else ""
            # In PDF, full-bleed slides would cover the body text. We keep the
            # cap at max_img_h (12cm) so the page still has room for narrative
            # text below. Treat slideBackground as a higher-priority imageFile.
            if slide_bg and not img_file:
                img_file = slide_bg

            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))
                story.append(Spacer(1, 0.3 * cm))

            # 섹션 이미지 — heading 다음, body 이전에 삽입 (fail-soft)
            # imageFile(사전 렌더) 우선, 없으면 imagePrompt(Bedrock) 사용
            img_abs = ""
            try:
                if img_file:
                    # 사전 렌더된 이미지 — Bedrock 호출 안 함
                    img_abs = img_file if os.path.isabs(img_file) else os.path.join(_local_root, img_file)
                    if not os.path.isfile(img_abs):
                        print(f"[generate_pdf] imageFile not found: {img_file}")
                        img_abs = ""
                elif img_prompt:
                    img_result_str = await _tool_generate_image(
                        {"prompt": img_prompt, "size": "1024x1024"},
                        project_path, aws_profile=aws_profile, bedrock_user=bedrock_user,
                    )
                    img_result = json.loads(img_result_str)
                    if "path" in img_result:
                        rel = img_result["path"]
                        cand = rel if os.path.isabs(rel) else os.path.join(_local_root, rel)
                        if os.path.isfile(cand):
                            img_abs = cand
            except Exception as e:
                # 이미지 단계 실패는 절대 outer try/except로 새지 않음 — 섹션은 텍스트 전용으로 진행
                print(f"[generate_pdf] image step failed section '{heading[:40]}': {e}")
                img_abs = ""

            if img_abs:
                try:
                    # 이미지의 실제 크기로 비율 보존하며 cap 적용
                    try:
                        from PIL import Image as _PIL
                        with _PIL.open(img_abs) as _im:
                            iw, ih = _im.size
                    except Exception:
                        iw, ih = 1024, 1024
                    scale = min(max_img_w / iw, max_img_h / ih, 1.0)
                    draw_w = iw * scale
                    draw_h = ih * scale
                    story.append(Image(img_abs, width=draw_w, height=draw_h))
                    story.append(Spacer(1, 0.3 * cm))
                except Exception as e:
                    print(f"[generate_pdf] embed failed section '{heading[:40]}': {e}")

            if body:
                for para in body.split("\n"):
                    if para.strip():
                        story.append(Paragraph(para, styles["Normal"]))
                        story.append(Spacer(1, 0.2 * cm))

        doc.build(story)
        # TASK 8 — build/save 후 디스크 자체 검증.
        if not os.path.isfile(output_path):
            return json.dumps({
                "error": "pdf-generation-failed",
                "detail": f"build reported success but file missing: {output_path}",
            })
        size_bytes = os.path.getsize(output_path)
        if size_bytes <= 0:
            try:
                os.remove(output_path)
            except OSError:
                pass
            return json.dumps({
                "error": "pdf-generation-failed",
                "detail": "build produced zero-byte file",
            })
        # Req 4.3: prefer reportlab's actual page counter; fall back to
        # section count if the attribute is unavailable.
        page_count = getattr(doc, "page", 0) or len(sections)
        return json.dumps({
            "path": relative_path,
            "absPath": output_path,  # TASK 8 근본수정 — 실제 저장 절대경로(카드 다운로드 정확성)
            "model": "reportlab",
            "pageCount": page_count,
            "sizeBytes": size_bytes,
            "fileSize": size_bytes,
        })
    except Exception as e:
        return json.dumps({"error": "pdf-generation-failed", "detail": str(e)[:200]})


def _open_presentation_with_timeout(path: str, timeout: int = 10):
    """Open a python-pptx Presentation from `path`, aborting if the open takes
    longer than `timeout` seconds (요구사항 9.1).

    python-pptx의 `Presentation()` 열기는 동기/블로킹이므로, 손상되거나 비정상적으로
    큰 파일이 무한정 멈추지 않도록 워커 스레드에서 실행하고
    `concurrent.futures`로 데드라인을 강제한다. 타임아웃 시
    `concurrent.futures.TimeoutError`(Exception 하위)를 던져 호출부의 폴백 경로가
    이를 잡도록 한다.
    """
    import concurrent.futures
    from pptx import Presentation

    executor = concurrent.futures.ThreadPoolExecutor(max_workers=1)
    try:
        future = executor.submit(Presentation, path)
        return future.result(timeout=timeout)
    finally:
        # 열기가 멈춰 있어도 프로세스 종료를 막지 않도록 워커 종료를 기다리지 않는다.
        executor.shutdown(wait=False)


def _clear_all_slides(prs) -> int:
    """템플릿에 들어있던 기존(샘플) 슬라이드를 모두 제거한다.

    사용자가 "템플릿을 활용해 작성"을 요청하면, 템플릿은 *스타일*(슬라이드 마스터/레이아웃/
    테마/색·폰트)로만 상속하고 내용은 새로 생성해야 한다. 그런데 템플릿 .pptx에 샘플
    슬라이드가 들어 있으면 그 위에 새 슬라이드가 *추가*되어, 기존 내용과 새 내용이 섞여
    "내용이 꼬이는" 버그가 발생한다. 이 함수는 슬라이드만 비우고 마스터/레이아웃/테마는
    그대로 보존한다(_resolve_layout / _safe_add_slide가 계속 동작).

    Returns:
        제거한 슬라이드 수.
    """
    removed = 0
    try:
        from pptx.oxml.ns import qn
    except Exception:
        qn = None
    try:
        sldIdLst = prs.slides._sldIdLst
        part = prs.part
        for sldId in list(sldIdLst):
            rId = None
            if qn is not None:
                try:
                    rId = sldId.get(qn("r:id"))
                except Exception:
                    rId = None
            # 관계(rel) 제거 — 슬라이드 파트 참조 해제
            if rId:
                try:
                    part.drop_rel(rId)
                except Exception:
                    pass
            # sldIdLst에서 항목 제거
            try:
                sldIdLst.remove(sldId)
                removed += 1
            except Exception:
                pass
    except Exception as e:
        print(f"[generate_pptx] 템플릿 기존 슬라이드 제거 실패(무시): {str(e)[:200]}")
    if removed:
        print(f"[generate_pptx] 템플릿 기존 슬라이드 {removed}개 제거 — 스타일만 상속하고 내용 새로 생성")
    return removed


def _clone_slide(prs, src):
    """템플릿의 디자인 슬라이드를 *복제*한다 — 배경/장식 도형/이미지를 모두 보존.

    Genspark 방식: 템플릿 슬라이드의 디자인(배경·도형·이미지)을 그대로 복사한 새 슬라이드를
    만든 뒤, 호출자가 텍스트만 새 내용으로 교체한다. 이미지 등 관계(rel)는 rId를 보존해
    deepcopy된 XML의 r:embed 참조가 깨지지 않게 한다.

    Returns:
        새 Slide 객체. 실패 시 None(호출자가 레이아웃 기반 슬라이드로 폴백).
    """
    import copy as _copy
    try:
        from pptx.oxml.ns import qn as _qn
        layout = src.slide_layout
        dst = prs.slides.add_slide(layout)
        # add_slide가 자동 생성한 placeholder 제거 — 디자인은 src에서 그대로 복사한다.
        for sp in list(dst.shapes):
            sp._element.getparent().remove(sp._element)
        # 배경(<p:bg>) 복사
        s_cSld = src._element.find(_qn('p:cSld'))
        d_cSld = dst._element.find(_qn('p:cSld'))
        if s_cSld is not None and d_cSld is not None:
            s_bg = s_cSld.find(_qn('p:bg'))
            if s_bg is not None:
                d_bg = d_cSld.find(_qn('p:bg'))
                if d_bg is not None:
                    d_cSld.remove(d_bg)
                d_cSld.insert(0, _copy.deepcopy(s_bg))
        # 도형 복사
        for sp in src.shapes:
            dst.shapes._spTree.append(_copy.deepcopy(sp._element))
        # 관계(이미지 등) 복사 — rId 보존(슬라이드 레이아웃 rel은 dst가 이미 보유하므로 제외)
        for rId, rel in src.part.rels.items():
            if rel.reltype.endswith('slideLayout'):
                continue
            if rId in dst.part.rels:
                continue
            try:
                if rel.is_external:
                    dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
                else:
                    dst.part.relate_to(rel.target_part, rel.reltype, rId)
            except Exception:
                pass
        return dst
    except Exception as e:
        print(f"[generate_pptx] 슬라이드 복제 실패(레이아웃 폴백): {str(e)[:200]}")
        return None


def _reset_designed_slide_text(slide):
    """디자인 슬라이드(템플릿 도너/복제본)의 *기존 텍스트만* 비운다.

    배경·이미지·장식 도형은 그대로 두고, 텍스트 프레임의 샘플 텍스트만 지워
    새 내용을 채울 빈 캔버스로 만든다. 어떤 예외도 밖으로 던지지 않는다.
    """
    try:
        shapes = list(slide.shapes)
    except Exception:
        return
    for sp in shapes:
        try:
            if sp.has_text_frame:
                sp.text_frame.clear()
        except Exception:
            continue


def _strip_slide_to_background(slide):
    """도너(템플릿) 슬라이드를 배경/로고만 남긴 깨끗한 캔버스로 정리한다.

    템플릿의 콘텐츠 placeholder/카드(텍스트박스, 빈 둥근사각형 등)는 우리가 생성하는
    내용(제목/다이어그램/불릿)과 겹쳐 "오와열이 깨지는" 원인이 되므로 제거한다.
    배경(<p:bg>), 그림(로고/장식 이미지), 텍스트 프레임이 없는 도형(라인/커넥터)은 보존하고,
    제목 placeholder는 우리 제목을 담기 위해 보존한다. 어떤 예외도 밖으로 던지지 않는다.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
    except Exception:
        _MST = None
    try:
        title_el = slide.shapes.title._element if slide.shapes.title is not None else None
    except Exception:
        title_el = None
    for sp in list(slide.shapes):
        try:
            if title_el is not None and sp._element is title_el:
                continue
            if _MST is not None and sp.shape_type == _MST.PICTURE:
                continue  # 로고/장식 이미지 보존
            if sp.has_text_frame:
                # 본문/콘텐츠 placeholder는 '그릇'으로 보존하되 샘플 텍스트만 비운다
                # → 텍스트 슬라이드 불릿이 누락되지 않는다(부분 누락 버그 수정).
                # 레이아웃 그리드에 정렬돼 있어 오와열도 유지된다.
                # 비-placeholder 장식 텍스트박스/카드만 제거(생성 콘텐츠와 충돌 방지).
                if bool(getattr(sp, "is_placeholder", False)):
                    try:
                        sp.text_frame.clear()
                    except Exception:
                        pass
                    continue
                sp._element.getparent().remove(sp._element)
        except Exception:
            continue


def _remove_content_band_pictures(slide, top_in: float, bottom_in: float) -> int:
    """네이티브 다이어그램 영역[top_in, bottom_in]을 가로지르는 얇고 넓은 템플릿
    장식 띠 PICTURE 를 제거한다(다이어그램과 겹쳐 보이는 문제, 이슈3).

    조건: PICTURE 이면서 폭 > 7in, 높이 < 1.2in, 세로 중심이 영역 안.
    로고/좌측 액센트바 등 작은 그림은 조건 불충족으로 보존된다.
    """
    removed = 0
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
    except Exception:
        return 0
    EMU = 914400.0
    for sp in list(slide.shapes):
        try:
            if sp.shape_type != _MST.PICTURE:
                continue
            w_in = (sp.width or 0) / EMU
            h_in = (sp.height or 0) / EMU
            t_in = (sp.top or 0) / EMU
            cy = t_in + h_in / 2.0
            if w_in > 7.0 and h_in < 1.2 and (top_in - 0.3) <= cy <= (bottom_in + 0.3):
                sp._element.getparent().remove(sp._element)
                removed += 1
        except Exception:
            continue
    return removed


def _remove_empty_text_shapes(slide):
    """텍스트가 비어 있는 텍스트박스/오토셰이프를 제거한다(제목·그림·라인 보존).

    네이티브 다이어그램이 본문을 대체하면서 비워진 텍스트박스, 또는 채우지 않은 콘텐츠
    카드가 빈 채로 남아 슬라이드가 지저분해지는 것을 막는다. 텍스트가 있는 도형은 보존한다.
    """
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
    except Exception:
        _MST = None
    try:
        title_el = slide.shapes.title._element if slide.shapes.title is not None else None
    except Exception:
        title_el = None
    for sp in list(slide.shapes):
        try:
            if title_el is not None and sp._element is title_el:
                continue
            if _MST is not None and sp.shape_type == _MST.PICTURE:
                continue
            # 네이티브 다이어그램의 장식 도형(카드/칩/액센트바=AUTO_SHAPE)은
            # 텍스트가 없어도 보존한다(텍스트는 별도 textbox). 이걸 제거하면
            # cards 가 배경 없이 텍스트만 떠 보인다(이슈4).
            if _MST is not None and sp.shape_type == _MST.AUTO_SHAPE:
                continue
            if not sp.has_text_frame:
                continue
            if (sp.text_frame.text or "").strip():
                continue  # 텍스트 있으면 보존
            sp._element.getparent().remove(sp._element)
        except Exception:
            continue


def _strip_text_over_fullbleed(slide):
    """B방향: 통짜(풀블리드) 배경 PICTURE가 채택된 슬라이드에서 그 배경 위에 겹치는
    편집 텍스트/도형(제목·부제 placeholder, 날짜 textbox, 네이티브 카드 AUTO_SHAPE,
    커넥터/라인 등)을 모두 제거해 '통짜 이미지 하나'만 남긴다(겹츨0·중복0).

    - 배경/로고 등 PICTURE 는 보존한다(통짜 배경 폐기 금지 — 손실0).
    - 배경(HTML 표지/본문 베이크 PNG 또는 Vertex 이미지)이 제목/본문 콘텐츠를
      포함한다고 신뢰하므로 네이티브 오버레이는 중복이며 제거 대상이다.
    - 어떤 예외도 밖으로 던지지 않는다. 반환: 제거한 shape 수(관측용)."""
    removed = 0
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
    except Exception:
        _MST = None
    for sp in list(slide.shapes):
        try:
            if _MST is not None and sp.shape_type == _MST.PICTURE:
                continue  # 통짜 배경/로고 그림 보존
            sp._element.getparent().remove(sp._element)
            removed += 1
        except Exception:
            continue
    return removed


def _clamp_shapes_into_bounds(slide):
    """경계 안전: 슬라이드의 모든 shape 를 (0,0,13.333,7.5)in 경계 안으로 후처리
    클램프한다. 음수 top/left 또는 경계 초과 shape 만 이동/축소하고, 이미 경계
    안이면 no-op(바이트 보존). 예외는 밖으로 던지지 않는다. 반환: 클램프한 shape 수."""
    EMU = 914400
    W = int(13.333 * EMU)
    H = int(7.5 * EMU)
    n = 0
    for sp in list(slide.shapes):
        try:
            l = int(sp.left) if sp.left is not None else 0
            t = int(sp.top) if sp.top is not None else 0
            w = int(sp.width) if sp.width is not None else 0
            h = int(sp.height) if sp.height is not None else 0
            nl, nt, nw, nh = l, t, w, h
            if nw > W:
                nw = W
            if nh > H:
                nh = H
            if nl < 0:
                nl = 0
            if nt < 0:
                nt = 0
            if nl + nw > W:
                nl = max(0, W - nw)
            if nt + nh > H:
                nt = max(0, H - nh)
            if (nl, nt, nw, nh) != (l, t, w, h):
                sp.left = nl
                sp.top = nt
                sp.width = nw
                sp.height = nh
                n += 1
        except Exception:
            continue
    return n


def _remove_slide_obj(prs, slide) -> bool:
    """프레젠테이션에서 특정 Slide 객체를 제거한다(sldIdLst 항목 + 관계).

    템플릿 도너 슬라이드 중 새 내용으로 사용되지 않은 잉여분을 제거할 때 사용.
    """
    try:
        from pptx.oxml.ns import qn as _qn
        target_part = slide.part
        rel_id = None
        for rid, rel in prs.part.rels.items():
            if rel.is_external:
                continue
            if rel.reltype.endswith('slideLayout') or rel.reltype.endswith('slideMaster'):
                continue
            if rel.reltype.endswith('slide') and rel.target_part is target_part:
                rel_id = rid
                break
        if rel_id is None:
            return False
        for sldId in list(prs.slides._sldIdLst):
            if sldId.get(_qn('r:id')) == rel_id:
                prs.slides._sldIdLst.remove(sldId)
                break
        try:
            prs.part.drop_rel(rel_id)
        except Exception:
            pass
        return True
    except Exception as e:
        print(f"[generate_pptx] 잉여 도너 슬라이드 제거 실패(무시): {str(e)[:160]}")
        return False


# 무템플릿 경로의 하드코딩 레이아웃 인덱스 (요구사항 5.2 하위 호환).
# 빈 Presentation()의 기본 레이아웃 순서: 0=Title Slide, 1=Title and Content,
# 3=Two Content. 무템플릿일 때는 이 매핑을 그대로 사용해 기존 동작을 보존한다.
LAYOUT_MAP = {"title": 0, "content": 1, "two-column": 3}


def _layout_has_content_placeholder(layout) -> bool:
    """레이아웃이 본문(body/object) placeholder(idx>=1)를 보유하면 True (요구사항 6.3).

    '콘텐츠 레이아웃' = 제목 외에 본문을 담을 수 있는 placeholder를 가진 레이아웃이다.
    enum import 실패를 포함해 어떤 예외에도 False를 반환하여 폴백 체인을 깨지 않는다.
    """
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        body_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
    except Exception:
        body_types = None
    try:
        for ph in layout.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf is None:
                continue
            idx = pf.idx
            if idx is None or idx < 1:
                continue
            if body_types is None:
                # enum 사용 불가 시 idx>=1 placeholder를 콘텐츠로 간주
                return True
            try:
                if pf.type in body_types:
                    return True
            except Exception:
                continue
    except Exception:
        return False
    return False


def _layout_name_matches(name_norm: str, layout_name: str) -> bool:
    """정규화된(소문자) 레이아웃 이름이 요청 layout_name 의미와 일치하는지 판정 (요구사항 6.2).

    - "title": 표지/제목 슬라이드/cover — '제목 및 내용'(title+content) 류는 제외
    - "content": 제목+내용 / content
    - "two-column": two content / comparison / 비교 / 2단
    """
    if layout_name == "title":
        # '제목 및 내용'(title+content)·'Content with Caption' 류를 표지로 오인하지 않도록 제외
        if "content" in name_norm or "내용" in name_norm:
            return False
        return (
            "title" in name_norm
            or "표지" in name_norm
            or "cover" in name_norm
            or "제목" in name_norm
        )
    if layout_name == "two-column":
        return (
            "two content" in name_norm
            or "two-content" in name_norm
            or "twocontent" in name_norm
            or "two column" in name_norm
            or "two-column" in name_norm
            or "comparison" in name_norm
            or "비교" in name_norm
            or "2단" in name_norm
        )
    # 기본 및 "content"
    return (
        "title and content" in name_norm
        or "제목 및 내용" in name_norm
        or "content" in name_norm
        or "내용" in name_norm
    )


def _resolve_layout(prs, layout_name: str, used_template: bool):
    """슬라이드 layout 이름을 실제 slide_layout 객체로 해석한다 (요구사항 5.2, 6.2, 6.3, 6.4).

    - used_template=False(무템플릿): 기존 LAYOUT_MAP 동작을 그대로 보존한다.
      prs.slide_layouts[LAYOUT_MAP.get(layout_name, 1)]을 반환하고, 인덱스가 범위를
      벗어나면 유효한 레이아웃으로 클램프한다 (하위 호환, 요구사항 5.2).
    - used_template=True(템플릿): 템플릿이 제공하는 레이아웃에 대해
      1) 이름 기반 매칭(요구사항 6.2) → 2) 첫 콘텐츠 레이아웃(body placeholder 보유, 요구사항 6.3)
      → 3) slide_layouts[0](요구사항 6.4) 폴백 체인을 적용한다.

    어떤 예외에도 raise하지 않고 slide_layouts[0]로 폴백한다.
    """
    try:
        if not used_template:
            # 무템플릿 — 기존 동작 보존 (요구사항 5.2)
            idx = LAYOUT_MAP.get(layout_name, 1)
            try:
                return prs.slide_layouts[idx]
            except IndexError:
                # 기존 per-slide 폴백 동작: slide_layouts[1]로 클램프
                try:
                    return prs.slide_layouts[1]
                except IndexError:
                    return prs.slide_layouts[0]

        layouts = list(prs.slide_layouts)

        # 1) 이름 기반 매칭 (요구사항 6.2)
        for layout in layouts:
            try:
                name_norm = (layout.name or "").strip().lower()
            except Exception:
                name_norm = ""
            if name_norm and _layout_name_matches(name_norm, layout_name):
                return layout

        # 2) 첫 콘텐츠 레이아웃 (요구사항 6.3)
        for layout in layouts:
            if _layout_has_content_placeholder(layout):
                return layout

        # 3) 콘텐츠 레이아웃 없음 → 인덱스 0 (요구사항 6.4)
        return prs.slide_layouts[0]
    except Exception:
        # 어떤 경우에도 raise 금지 — 마지막 안전망
        return prs.slide_layouts[0]


def _safe_add_slide(prs, layout):
    """slide_layout으로 슬라이드를 추가한다. 템플릿 레이아웃 편차로 add_slide가 실패하면
    (요구사항 9.5) 안전한 폴백 레이아웃(slide_layouts[0] → 사용 가능한 임의 레이아웃)으로
    재시도한다. 무템플릿 happy-path에서는 prs.slides.add_slide(layout)와 동일하게 동작해
    기존 출력을 그대로 보존한다(요구사항 5.2)."""
    try:
        return prs.slides.add_slide(layout)
    except Exception as e:
        # 요구사항 9.5 — 템플릿 레이아웃 적용 실패를 격리하고 안전한 기본 레이아웃으로 재시도
        print(f"[generate_pptx] add_slide 실패 → 폴백 레이아웃 재시도: {str(e)[:200]}")
        try:
            return prs.slides.add_slide(prs.slide_layouts[0])
        except Exception:
            for _lay in prs.slide_layouts:
                try:
                    return prs.slides.add_slide(_lay)
                except Exception:
                    continue
            raise


def _safe_set_title(slide, text):
    """슬라이드 제목을 편집 가능한 텍스트로 채운다(요구사항 6.8). 템플릿 레이아웃에 title
    placeholder가 없거나 설정이 실패하면(요구사항 9.5) 텍스트박스로 폴백해 제목 콘텐츠를
    보존하며, 어떤 예외도 밖으로 던지지 않는다. 무템플릿 happy-path에서는
    `slide.shapes.title.text = text`와 동일하게 동작해 기존 출력을 보존한다(요구사항 5.2)."""
    try:
        title_ph = slide.shapes.title
    except Exception:
        title_ph = None
    if title_ph is not None:
        try:
            title_ph.text = text
            return
        except Exception as e:
            print(f"[generate_pptx] title placeholder 설정 실패 → textbox 폴백: {str(e)[:200]}")
    # 폴백: 텍스트박스로 제목 콘텐츠 유지 (요구사항 6.8, 9.5) — 템플릿 단계 실패가 콘텐츠 손실로
    # 이어지지 않도록 한다(요구사항 9.6).
    try:
        from pptx.util import Inches as _In
        tb = slide.shapes.add_textbox(_In(0.6), _In(0.3), _In(12.1), _In(1.0))
        tb.text_frame.text = text
    except Exception as e:
        print(f"[generate_pptx] title textbox 폴백 실패: {str(e)[:200]}")


def _remove_empty_placeholders(slide):
    """빈 placeholder를 슬라이드에서 *삭제*한다.

    PowerPoint는 비어 있는 본문/콘텐츠 placeholder를 "텍스트를 입력하십시오" 프롬프트
    텍스트 + 점선 테두리로 렌더한다. 우리가 채우지 않은(네이티브 다이어그램으로 대체했거나
    bullets가 없는) placeholder를 그대로 두면 슬라이드가 "중구난방"으로 보인다.
    `text_frame.clear()`만으로는 프롬프트가 사라지지 않으므로 shape 요소 자체를 제거한다.

    - 제목(title/center-title) placeholder는 항상 보존.
    - 텍스트가 있는 placeholder는 보존(bullets 채운 본문 등).
    - 그림(picture) placeholder는 건드리지 않음.
    어떤 예외도 밖으로 던지지 않는다(템플릿 편차 격리, 요구사항 9.5).
    """
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER as _PP
    except Exception:
        _PP = None
    try:
        placeholders = list(slide.placeholders)
    except Exception:
        return
    for ph in placeholders:
        try:
            fmt = ph.placeholder_format
            ptype = fmt.type
            idx = fmt.idx
        except Exception:
            continue
        # 제목류 보존
        is_title = False
        if _PP is not None:
            try:
                is_title = ptype in (_PP.TITLE, _PP.CENTER_TITLE)
            except Exception:
                is_title = (idx == 0)
        else:
            is_title = (idx == 0)
        if is_title:
            continue
        # 그림 placeholder는 보존(이미지가 들어갈 수 있음)
        if _PP is not None:
            try:
                if ptype == _PP.PICTURE:
                    continue
            except Exception:
                pass
        # 텍스트가 있으면 보존
        has_text = False
        try:
            has_text = bool((ph.text_frame.text or "").strip())
        except Exception:
            has_text = False
        if has_text:
            continue
        # 빈 placeholder → shape 요소 삭제 (프롬프트/점선 제거)
        try:
            ph._element.getparent().remove(ph._element)
        except Exception:
            pass


def _wants_high_quality_slides(*texts) -> bool:
    """사용자 의도에서 '고품질/PDF용/젠스파크급' 신호 감지 → HTML 렌더 경로(하이브리드).

    기본은 편집 가능 네이티브. 아래 키워드가 잡힌 요청에 한해 HTML(Chromium 렌더)
    풀블리드 고품질 경로를 켠다. 브리지(Electron) 미가용 시 자동 네이티브 폴백.
    """
    blob = " ".join(str(t or "") for t in texts).lower()
    if not blob.strip():
        return False
    kw = (
        "고품질", "고퀄", "젠스파크", "genspark", "gamma", "감마",
        "pdf용", "발표자료", "프레젠테이션급", "프리미엄", "premium",
        "픽셀", "pixel", "동급", "고급 디자인", "high quality", "high-quality",
        "정교한 디자인", "디자인 퀄리티",
    )
    return any(k in blob for k in kw)


def _wants_editable_slides(*texts) -> bool:
    """사용자가 '편집 가능/네이티브 도형' 결과물을 *명시적으로* 요청했는지 감지(옵트아웃).

    품질은 기본 고품질(HTML 렌더)이고, 아래 신호가 있을 때만 편집 가능 네이티브로 전환한다.
    오탐을 줄이기 위해 단어 'edit/편집' 단독이 아니라 명시적 구절만 매칭한다.
    """
    blob = " ".join(str(t or "") for t in texts).lower()
    if not blob.strip():
        return False
    kw = (
        "편집 가능", "편집가능", "편집이 가능", "편집할 수 있", "수정 가능", "수정가능",
        "수정할 수 있", "고칠 수 있", "네이티브", "도형으로", "도형 편집", "개체 편집",
        "ppt에서 편집", "ppt에서 수정", "powerpoint에서 편집", "파워포인트에서 편집",
        "editable", "edit in powerpoint", "native shape",
    )
    return any(k in blob for k in kw)


def _resolve_html_slides(env_val, *intent_texts, tool_highquality=False) -> bool:
    """HTML 고품질 렌더 사용 여부 결정 (기본 ON, 편집 가능 요청 시 OFF).

    - 환경변수 AE_ENABLE_HTML_SLIDES == "1" → 항상 ON (강제)
    - == "0" → 항상 OFF (강제, 편집 가능 네이티브)
    - 그 외(미설정) → 기본 ON. 단 사용자가 '편집 가능/네이티브'를 명시하면 OFF.
    tool_highquality=True(도구 입력 highQuality)면 편집 요청보다 우선해 ON.
    """
    env = (env_val or "").strip()
    if env == "1":
        return True
    if env == "0":
        return False
    if tool_highquality:
        return True
    blob = " ".join(str(t or "") for t in intent_texts).lower()
    # 기본 ON — Genspark/Gamma급 HTML 풀블리드가 기본 고품질 산출물이다.
    # 'PDF용 고정 이미지/편집 불필요' 명시는 ON을 굳히고, '편집 가능/네이티브/수정 가능'
    # 명시일 때만 OFF(편집 가능 네이티브)로 전환한다. (HTML-positive 키워드를 먼저
    # 검사해 'non-editable'이 editable로 오탐되지 않게 한다.)
    _html_kw = ("pdf용", "pdf 용", "이미지로 고정", "고정 이미지",
                "편집 불필요", "flat image", "non-editable")
    if any(k in blob for k in _html_kw):
        return True
    # 기본 ON — 사용자가 '편집 가능/네이티브/수정 가능'을 명시한 경우에만 OFF.
    _native_kw = ("편집 가능", "편집가능", "수정 가능", "수정가능",
                  "네이티브", "editable", "native")
    if any(k in blob for k in _native_kw):
        return False
    return True


def _apply_universal_slide_design(prs, palette=None, doc_title="", used_template=False, skip_slide_ids=None):
    """무템플릿 PPTX의 모든 슬라이드에 공통 디자인을 적용(내용/다이어그램 무관).

    상단 액센트 바 · 제목 색/폰트 통일 · 제목 밑줄 액센트 · 본문 글머리표 폰트/색/간격
    · 하단 푸터(문서명 + 페이지). 표지(첫 슬라이드)는 더 큰 강조 밴드 + 중앙 정렬.
    어떤 예외도 밖으로 던지지 않는다(미화 실패가 생성 실패로 이어지지 않게).
    """
    try:
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except Exception:
        return

    def _hx(h):
        s = (h or "").lstrip("#")
        try:
            return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
        except Exception:
            return RGBColor(0x44, 0x72, 0xC4)

    pal = []
    if isinstance(palette, (list, tuple)):
        for c in palette:
            if isinstance(c, str) and len(c.lstrip("#")) == 6:
                pal.append(c.lstrip("#"))
    primary = pal[0] if len(pal) >= 1 else "4472C4"
    accent = pal[1] if len(pal) >= 2 else "ED7D31"
    MUTED = RGBColor(0x8A, 0x8F, 0x98)
    DARK = RGBColor(0x33, 0x36, 0x3D)

    try:
        W = prs.slide_width / 914400.0
        H = prs.slide_height / 914400.0
    except Exception:
        return

    def _norect(shp):
        try:
            shp.line.fill.background()
        except Exception:
            pass

    def _font(run, size=None, bold=None, color=None):
        try:
            run.font.name = "Malgun Gothic"
        except Exception:
            pass
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold
        if color is not None:
            run.font.color.rgb = color

    slides = list(prs.slides)
    total = len(slides)
    # task 8.5 — 하이브리드 content 편집 경로(render_native_layout)로 렌더된 슬라이드는
    # 이미 자족적 chrome(제목 TextBox·푸터·번호)을 가진다. 여기서 덱 레벨 chrome(헤더
    # 밴드·번호 배지·푸터)을 또 얹으면 제목/푸터가 이중 방출돼 겹친다(Property 18 위반).
    # 해당 슬라이드는 id 로 식별해 건너뛴다 — 단일 chrome 소스 유지. skip 집합에 없는
    # 비하이브리드/비-content 슬라이드는 바이트 동일(가산적·게이팅).
    _skip_ids = skip_slide_ids if isinstance(skip_slide_ids, (set, frozenset)) else set()
    for i, slide in enumerate(slides):
        is_cover = (i == 0)
        if is_cover and used_template:
            continue
        if id(slide) in _skip_ids:
            continue
        try:
            # HTML 풀블리드 배경이 이미 깔린 슬라이드는 장식을 얹지 않는다(겹침 방지).
            _has_fullbleed = False
            try:
                from pptx.enum.shapes import MSO_SHAPE_TYPE as _MST
                _sw_emu = int(prs.slide_width * 0.95)
                _sh_emu = int(prs.slide_height * 0.95)
                for _sh in slide.shapes:
                    if _sh.shape_type == _MST.PICTURE and (_sh.width or 0) >= _sw_emu and (_sh.height or 0) >= _sh_emu:
                        _has_fullbleed = True
                        break
            except Exception:
                _has_fullbleed = False
            if _has_fullbleed:
                continue
            # ── 헤더 ──
            title_shp = None
            try:
                title_shp = slide.shapes.title
            except Exception:
                title_shp = None
            if is_cover:
                bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                             Inches(0), Inches(0), Inches(W), Inches(0.20))
                bar.fill.solid(); bar.fill.fore_color.rgb = _hx(primary); _norect(bar)
                tick = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                              Inches(0), Inches(0), Inches(2.0), Inches(0.20))
                tick.fill.solid(); tick.fill.fore_color.rgb = _hx(accent); _norect(tick)
                if title_shp is not None:
                    try:
                        for p in title_shp.text_frame.paragraphs:
                            p.alignment = PP_ALIGN.CENTER
                            runs = p.runs or []
                            if not runs and (p.text or ""):
                                r0 = p.add_run(); r0.text = p.text; runs = [r0]
                            for r in runs:
                                _font(r, size=40, bold=True, color=_hx(primary))
                    except Exception:
                        pass
                    try:
                        tl = title_shp.left / 914400.0; tt = title_shp.top / 914400.0
                        tw = title_shp.width / 914400.0; th = title_shp.height / 914400.0
                        uw = min(2.2, max(1.2, tw * 0.22)); ux = tl + (tw - uw) / 2.0; uy = tt + th + 0.04
                        if uy < H - 0.6:
                            ul = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                        Inches(ux), Inches(uy), Inches(uw), Inches(0.07))
                            try:
                                ul.adjustments[0] = 0.5
                            except Exception:
                                pass
                            ul.fill.solid(); ul.fill.fore_color.rgb = _hx(accent); _norect(ul)
                    except Exception:
                        pass
            else:
                # 본문: 업로드 PDF급 다크 헤더 밴드 + 번호 배지 + 흰 제목.
                _sec_title = ""
                try:
                    if title_shp is not None and title_shp.has_text_frame:
                        _sec_title = (title_shp.text_frame.text or "").strip()
                        title_shp.text_frame.clear()
                        try:
                            title_shp.top = Emu(int(0.02 * 914400))
                            title_shp.height = Emu(int(0.30 * 914400))
                        except Exception:
                            pass
                except Exception:
                    pass
                band_h = 1.0
                try:
                    band = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                  Inches(0.0), Inches(0.0), Inches(W), Inches(band_h))
                    try:
                        band.adjustments[0] = 0.04
                    except Exception:
                        pass
                    band.fill.solid(); band.fill.fore_color.rgb = RGBColor(0x2B, 0x38, 0x46); _norect(band)
                except Exception:
                    pass
                try:
                    strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                   Inches(0.0), Inches(0.0), Inches(0.16), Inches(band_h))
                    strip.fill.solid(); strip.fill.fore_color.rgb = _hx(accent); _norect(strip)
                except Exception:
                    pass
                try:
                    bd = 0.5; bx = 0.55; by = (band_h - bd) / 2.0
                    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                   Inches(bx), Inches(by), Inches(bd), Inches(bd))
                    try:
                        badge.adjustments[0] = 0.25
                    except Exception:
                        pass
                    badge.fill.solid(); badge.fill.fore_color.rgb = _hx(primary); _norect(badge)
                    _btf = badge.text_frame
                    try:
                        _btf.margin_left = Emu(0); _btf.margin_right = Emu(0)
                        _btf.margin_top = Emu(0); _btf.margin_bottom = Emu(0)
                    except Exception:
                        pass
                    _bp = _btf.paragraphs[0]; _bp.alignment = PP_ALIGN.CENTER
                    _br = _bp.add_run(); _br.text = str(i).zfill(2)
                    _font(_br, size=15, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
                except Exception:
                    pass
                try:
                    from pptx.enum.text import MSO_ANCHOR as _MA2
                    tbx = slide.shapes.add_textbox(Inches(1.25), Inches(0.0), Inches(W - 2.5), Inches(band_h))
                    _tf2 = tbx.text_frame; _tf2.word_wrap = True
                    try:
                        _tf2.vertical_anchor = _MA2.MIDDLE
                    except Exception:
                        pass
                    _tp = _tf2.paragraphs[0]; _tp.alignment = PP_ALIGN.LEFT
                    _tr = _tp.add_run(); _tr.text = _sec_title[:60]
                    _font(_tr, size=24, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
                except Exception:
                    pass

            # ── 본문 글머리표 스타일 통일 ──
            try:
                from pptx.enum.shapes import PP_PLACEHOLDER as _PP
            except Exception:
                _PP = None
            for ph in (list(slide.placeholders) if hasattr(slide, "placeholders") else []):
                try:
                    if title_shp is not None and ph._element is title_shp._element:
                        continue
                    fmt = ph.placeholder_format
                    if _PP is not None and fmt.type in (_PP.TITLE, _PP.CENTER_TITLE):
                        continue
                    if not ph.has_text_frame:
                        continue
                    if not (ph.text_frame.text or "").strip():
                        continue
                    for p in ph.text_frame.paragraphs:
                        try:
                            p.space_after = Pt(8)
                        except Exception:
                            pass
                        runs = p.runs or []
                        for r in runs:
                            _font(r, size=16, color=DARK)
                except Exception:
                    continue

            # ── 하단 푸터(문서명 + 페이지) — 표지 제외 ──
            if not is_cover:
                try:
                    hair = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                                  Inches(0.6), Inches(H - 0.42),
                                                  Inches(W - 1.2), Inches(0.012))
                    hair.fill.solid()
                    hair.fill.fore_color.rgb = RGBColor(0xDD, 0xDF, 0xE4)
                    _norect(hair)
                except Exception:
                    pass
                try:
                    fl = slide.shapes.add_textbox(Inches(0.6), Inches(H - 0.40),
                                                  Inches(W - 2.0), Inches(0.32))
                    fp = fl.text_frame.paragraphs[0]
                    fr = fp.add_run(); fr.text = str(doc_title or "")[:70]
                    _font(fr, size=9, color=MUTED)
                except Exception:
                    pass
                try:
                    pn = slide.shapes.add_textbox(Inches(W - 1.4), Inches(H - 0.40),
                                                  Inches(0.8), Inches(0.32))
                    pp = pn.text_frame.paragraphs[0]
                    pp.alignment = PP_ALIGN.RIGHT
                    pr = pp.add_run(); pr.text = f"{i + 1} / {total}"
                    _font(pr, size=9, color=MUTED)
                except Exception:
                    pass
        except Exception:
            continue


def _add_md_paragraph(p, text):
    """마크다운 볼드(**...**)를 굵게 렌더하고 나머지 마커는 제거해
    PPTX 단락을 run으로 채운다. raw 마크다운 노출 문제 해결."""
    import re as _re_md
    raw = str(text or "")
    raw = _re_md.sub(r"`([^`]+)`", r"\1", raw)            # 인라인 코드 백틱 제거
    raw = _re_md.sub(r"^\s*(?:[-*•]\s+|#{1,6}\s+)", "", raw)  # 선두 글머표/헤더 제거
    parts = _re_md.split(r"(\*\*[^*]+\*\*)", raw)
    made = False
    for seg in parts:
        if not seg:
            continue
        r = p.add_run()
        if seg.startswith("**") and seg.endswith("**") and len(seg) > 4:
            r.text = seg[2:-2]
            r.font.bold = True
        else:
            r.text = seg.replace("**", "")
        made = True
    if not made:
        p.add_run().text = ""


def _should_native_render(sd: dict, layout: str, html_enabled: bool) -> bool:
    """본문 콘텐츠 슬\�라이드를 네이티브_렌더러로 라우팅할지 결정 (design §Components §2, Property 3).

    True 조건 (모두 충족 시):
      - ``html_enabled`` — HTML 고밀도 경로가 활성(네이티브 라우팅은 콘텐츠 베이크
        HTML→PNG 경로를 대체하는 분기다). html_enabled 가 False 면 기존
        ``AE_PREFER_EDITABLE_DIAGRAM`` 네이티브 다이어그램 경로를 보존한다(Req 9.3).
      - ``layout`` in ``NATIVE_LAYOUT_REGISTRY`` — 알려진_레이아웃 한정(Non-Goals).
      - caller 가 ``imageFile``/``slideBackground``/``nativeDiagram`` 을 명시하지
        않음 — 명시 우선순위 보존(Req 9.1, Property 3).
      - 슬라이드에 콘텐츠 텍스트(제목 또는 불릿)가 존재.

    그 외에는 False 를 반환해 기존 명시/베이크 경로를 주 렌더러로 보존한다(additive,
    no-op). LLM/게이트웨이/네트워크 호출이 없는 순수 결정 함수다(결정성).
    """
    try:
        if not html_enabled:
            return False
        if not isinstance(sd, dict):
            return False
        # caller 명시 우선순위 보존 (Req 9.1, Property 3)
        if sd.get("imageFile") or sd.get("slideBackground") or sd.get("nativeDiagram"):
            return False
        # 알려진_레이아웃만 (cover/section_divider/two_column/feature_grid/
        # timeline/comparison/architecture). 비-네이티브 레이아웃은 기존 베이크 보존.
        try:
            from ai_engine.native_layout_renderer import NATIVE_LAYOUT_REGISTRY as _NLR
        except ImportError:
            from native_layout_renderer import NATIVE_LAYOUT_REGISTRY as _NLR
        if not layout or layout not in _NLR:
            return False
        # 콘텐츠 텍스트 존재 여부
        _title = str(sd.get("title", "") or "").strip()
        _bullets = sd.get("bullets") or []
        if isinstance(_bullets, (list, tuple)):
            _has_bullets = any(str(b).strip() for b in _bullets)
        else:
            _has_bullets = bool(str(_bullets).strip())
        return bool(_title or _has_bullets)
    except Exception:
        return False



async def _tool_generate_pptx(tool_input: dict, project_path: str, aws_profile: str = '', bedrock_user: str = '') -> str:  # [patched-credentials]
    """Generate a PowerPoint presentation using python-pptx. Accepts lenient input shapes."""
    title, slides_data = _normalize_doc_input(tool_input, default_kind="slides")

    if not title:
        return json.dumps({"error": "title is required"})
    if not slides_data:
        return json.dumps({"error": "slides is required"})

    try:
        from pptx import Presentation
        from pptx.util import Inches, Emu
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "python-pptx", "hint": "pip install python-pptx"})

    import time as _t, re as _re
    # Always use a locally-existing directory for generated media.
    # Remote project_path may not exist locally; fall back to cwd.
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _slug_from_title(title) or "deck"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.pptx"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    # 요구사항 5.2/6.1/6.9/9.1 — 템플릿 분기. 호출부가 해석해 전달한 templatePath가
    # 있으면 그 기준 .pptx를 열어 마스터/레이아웃/테마를 상속한다. 없으면 기존과
    # 바이트 동일한 무템플릿 경로(빈 Presentation + 16:9 리사이즈)로 진행한다.
    template_path = tool_input.get("templatePath", "") if isinstance(tool_input, dict) else ""
    template_id = tool_input.get("templateId", "") if isinstance(tool_input, dict) else ""
    used_template = False

    # 요구사항 9.2 — tool_input에 styleProfile이 함께 전달될 수 있다. 손상/형식 불일치
    # styleProfile이 _tool_generate_pptx를 중단시키지 않도록, 읽기·역직렬화를 격리해 실패 시
    # 원인을 로그(≤200자)로 남기고 기본값(None)으로 폴백한다. 토큰 적용 자체는 후속 단계
    # (slide_templates/배경 파이프라인)의 책임이며, 여기서는 절대 raise하지 않음을 보장한다.
    # 무템플릿 출력에는 영향을 주지 않는다(요구사항 5.2 바이트 동일 보존).
    style_profile = None
    if isinstance(tool_input, dict) and tool_input.get("styleProfile") is not None:
        try:
            _sp_raw = tool_input.get("styleProfile")
            if isinstance(_sp_raw, dict):
                style_profile = _sp_raw
            elif isinstance(_sp_raw, str) and _sp_raw.strip():
                style_profile = json.loads(_sp_raw)
                if not isinstance(style_profile, dict):
                    raise ValueError("styleProfile is not a JSON object")
        except Exception as e:
            # 요구사항 9.2 — 손상 styleProfile은 SLIDE_DESIGN 기본값(None)으로 폴백
            print(f"[generate_pptx] styleProfile load 실패 → 기본값 폴백: {str(e)[:200]}")
            style_profile = None

    # 문제2 — 네이티브 다이어그램용 색 팔레트. styleProfile이 있으면 거기서 파생,
    # 없으면 None(빌더가 기본 색 사용). _build_palette는 어떤 입력에도 raise 안 함.
    try:
        _tpl_palette_for_native = _build_palette(style_profile)
    except Exception:
        _tpl_palette_for_native = None

    # === Vertex 고품질 배경(직접 호출 경로) — 표지 hero + 본문 공유 배경 ============
    # 모델이 generate_pptx를 직접 호출하는 경로에서도 force-generate 경로와 동일하게
    # 표지/본문에 Vertex 고품질 배경을 깐다. 그 위에는 편집 가능한 네이티브 표지/
    # 다이어그램/텍스트가 올라간다. AE_DISABLE_VERTEX_HERO / AE_DISABLE_VERTEX_BODY_BG로 끔.
    _dp_body_bg = ""
    try:
        import asyncio as _dp_aio
        _dp_vertex_on = False
        # Vertex 자체 AI 배경은 기본 OFF — 템플릿 스타일을 우선한다. 옵트인 시만.
        if (os.environ.get("AE_ENABLE_VERTEX_BG", "").strip() == "1"
                and os.environ.get("AE_DISABLE_VERTEX_IMAGE", "").strip() != "1"):
            try:
                from ai_engine.vertex_image_module import get_vertex_image_client as _dp_vg
            except Exception:
                try:
                    from vertex_image_module import get_vertex_image_client as _dp_vg
                except Exception:
                    _dp_vg = None
            if _dp_vg is not None:
                try:
                    _dp_vertex_on = bool(_dp_vg(aws_profile=aws_profile or "").enabled)
                except Exception:
                    _dp_vertex_on = False
        if _dp_vertex_on:
            _dp_title = str(title or "presentation")[:90]
            _cb_now = tool_input.get("coverBackground", "") if isinstance(tool_input, dict) else ""
            _any_bg = any(isinstance(s, dict) and s.get("slideBackground")
                          for s in (slides_data or []))
            _disable_all = (os.environ.get("AE_DISABLE_VERTEX_HERO", "").strip() == "1"
                            and os.environ.get("AE_DISABLE_VERTEX_BODY_BG", "").strip() == "1")
            if not _disable_all:
                _bgp = ("Professional abstract background for a corporate presentation about: "
                        + _dp_title + ". Soft smooth gradient with subtle geometric depth, "
                        "calm muted tones, generous clean negative space, cinematic soft lighting, "
                        "no text, no words, no letters, no charts, high resolution, 16:9.")
                try:
                    _bgr = await _try_vertex_image_single(_bgp, "1280x720", project_path, aws_profile)
                except Exception as _be:
                    _bgr = None
                    print(f"[generate_pptx] Vertex 공유 배경 실패: {str(_be)[:140]}")
                if isinstance(_bgr, dict) and _bgr.get("path"):
                    _shared = _bgr["path"]
                    if (not _cb_now and isinstance(tool_input, dict)
                            and os.environ.get("AE_DISABLE_VERTEX_HERO", "").strip() != "1"):
                        tool_input["coverBackground"] = _shared
                    if (not _any_bg
                            and os.environ.get("AE_DISABLE_VERTEX_BODY_BG", "").strip() != "1"):
                        _dp_body_bg = _shared
                    print(f"[generate_pptx] Vertex 공유 배경 적용: {_shared}")
    except Exception as _dp_e:
        print(f"[generate_pptx] Vertex 배경 단계 예외(무시): {str(_dp_e)[:160]}")

    try:
        prs = None
        if template_path:
            try:
                # 요구사항 6.1, 9.1 — 10초 타임아웃 + 예외 격리
                prs = _open_presentation_with_timeout(template_path, timeout=10)
                used_template = True
                # 템플릿의 디자인 슬라이드(배경·장식·이미지)를 도너로 재사용한다.
                # (Genspark 방식: 디자인을 복사해 새 내용을 채운다. 슬라이드를 비우지 않는다.)
            except Exception as e:
                # 요구사항 6.9, 9.1 — 파일 없음·손상·형식 불일치·타임아웃 시 무템플릿 폴백
                print(f"[generate_pptx] template open failed → no-template fallback: {str(e)[:200]}")
                prs = None
                used_template = False
        if prs is None:
            # baseline 경로 (요구사항 5.2) — templateId/templatePath 부재 시 기존과 동일
            prs = Presentation()
            # Force 16:9 widescreen — Genspark/Gamma-class slides are 1920×1080.
            # python-pptx default is 10×7.5 in (4:3); we need 13.333×7.5 in to
            # fit a 1920×1080 PNG full-bleed without letterboxing.
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

        # === 템플릿 디자인 슬라이드 재사용(Genspark 방식) 셋업 ===
        # 템플릿이면 그 안의 디자인 슬라이드를 "도너"로 잡아, 새 슬라이드마다 도너 디자인
        # (배경·장식·이미지)을 재사용하고 텍스트만 새 내용으로 교체한다. 도너보다 슬라이드가
        # 많이 필요하면 콘텐츠 도너를 복제한다. 무템플릿이면 기존처럼 레이아웃에서 새로 만든다.
        _donors = []
        if used_template:
            try:
                _donors = list(prs.slides)
            except Exception:
                _donors = []
        _use_designed = len(_donors) > 0
        # 콘텐츠용 도너: 본문 placeholder를 가진 도너 우선(없으면 2번째, 그것도 없으면 첫째).
        _content_donor = None
        if _use_designed:
            for _d in _donors:
                try:
                    if len(_d.placeholders) > 1:
                        _content_donor = _d
                        break
                except Exception:
                    continue
            if _content_donor is None:
                _content_donor = _donors[1] if len(_donors) >= 2 else _donors[0]
        _donor_consumed = {"n": 0}

        # === Genspark/Gamma급 HTML 디자인 슬라이드 (기본 ON) ===
        # 무템플릿이고 Electron 브리지(render-html-to-png)가 가용하면 각 슬라이드를
        # 고품질 HTML 레이아웃으로 렌더해 풀블리드 배경으로 사용한다. 브리지가 없으면
        # (헤드리스/테스트) 자동으로 네이티브 도형 경로로 폴백한다. 템플릿 사용 시에는
        # 템플릿 디자인을 유지하기 위해 HTML 오버라이드를 적용하지 않는다.
        # 기본 고품질(HTML 렌더) — '대충/넓게' 요청해도 고품질로 나온다.
        # 사용자가 '편집 가능/네이티브'를 명시하면 OFF(옵트아웃). 환경변수로 강제 가능.
        # 레거시 AE_DISABLE_HTML_SLIDES=1 은 계속 강제 OFF로 존중.
        _html_on = (
            _resolve_html_slides(
                os.environ.get("AE_ENABLE_HTML_SLIDES", ""),
                title,
                tool_highquality=bool(isinstance(tool_input, dict) and tool_input.get("highQuality")),
            )
            and os.environ.get("AE_DISABLE_HTML_SLIDES", "").strip() != "1"
        )
        _html_enabled = False
        _html_gw = None
        _html_model = ""
        # === renderReport 관측 상태 (additive — 제어흐름/출력 바이트 불변) ===
        # 선택된 렌더 경로를 도구 결과에 노출해, 덱이 실제로 HTML 풀블리드/Vertex
        # 경로를 썼는지(또는 왜 폴백했는지)를 호출자가 알 수 있게 한다. 값은 기존
        # 결정 지점(_html_on/_bridge_ok/_chrome_ok, 아래 Vertex 게이트)을 그대로
        # 재사용해 채운다(추가 LLM/Vertex/브리지 호출 없음).
        _bridge_ok = False
        _chrome_ok = False
        _rr_html_renderer = ""
        _rr_html_reason = ""
        _rr_vertex_enabled = False
        _rr_vertex_reason = ""
        _rr_cover_path = ""
        _rr_slides = []
        # HTML 고품질 렌더는 무템플릿/템플릿 모두에서 사용한다.
        # 템플릿이 활성이면 그 Style_Profile(색/폰트)을 HTML 디자인 토큰에 주입해
        # "젠스파크급 레이아웃 + 템플릿 브랜드 색"을 동시에 얻는다(아래 _generate_html_slide_for_section
        # 호출이 style_profile을 전달). 과거: used_template이면 HTML을 꺼서 항상 휑한 네이티브
        # 도너로 빠졌음(진단 로그로 확정된 근본 원인).
        if _html_on:
            try:
                _bridge_ok = _call_bridge("status", {}) is not None
                _chrome_ok = bool(_find_local_chrome())
                if _bridge_ok or _chrome_ok:
                    _html_gw = _get_gw(aws_profile, bedrock_user)
                    _html_model = _specialized_model_for_task(
                        "file_generation", "", aws_profile=aws_profile, bedrock_user=bedrock_user)
                    _html_enabled = True
                    _renderer = "Electron 브리지" if _bridge_ok else "로컬 Chrome"
                    print(f"[generate_pptx] HTML 디자인 슬라이드 활성 (Genspark급 풀블리드, 렌더러={_renderer}, 템플릿={used_template})")
            except Exception as _he0:
                print(f"[generate_pptx] HTML 슬라이드 비활성(브리지 확인 실패): {str(_he0)[:160]}")
                _html_enabled = False

        # === 하이브리드 렌더 Feature Flag 게이트 (pptx-ultra-quality-hybrid-render) ===
        # `_html_enabled` 결정 직후, opt-in 플래그 `AE_HYBRID_RENDER`를 정확히 1회 읽는다.
        # 파싱은 순수 함수 `_hybrid_render_enabled`에 위임(결정론적, raise 없음).
        #   "1" → True / 미설정·""·"0" → False / 인식 불가 값 → False + 경고 1줄
        # `_hybrid_on == False`이면 이후 하이브리드 분기(task 7.2~7.4)는 전부 no-op이 되어
        # 기존 제어흐름·출력이 바이트 단위로 동일하게 유지된다(R6.1/6.4). 본 태스크(7.1)는
        # 게이트 변수만 도입하며 라우팅 로직은 추가하지 않는다.
        _hybrid_on = _hybrid_render_enabled(os.environ.get("AE_HYBRID_RENDER", ""))

        # === 진단 로그 (AE_GEN_DIAG) — HTML 분기 결정 추적 ===
        try:
            import datetime as _dtg
            with open("/tmp/ae_gen_diag.log", "a", encoding="utf-8") as _dg:
                _dg.write(
                    f"{_dtg.datetime.now().isoformat()} generate_pptx "
                    f"title={str(title)[:40]!r} used_template={used_template} "
                    f"_html_on={_html_on} _html_enabled={_html_enabled} "
                    f"chrome={bool(_find_local_chrome())} "
                    f"bridge={_call_bridge('status', {}) is not None} "
                    f"slides={len(slides_data) if isinstance(slides_data, list) else '?'}\n"
                )
        except Exception:
            pass

        # === renderReport: HTML 경로 판정 (관측만) ===
        if _html_enabled:
            _rr_html_renderer = "bridge" if _bridge_ok else "local-chrome"
            _rr_html_reason = ""
        else:
            _rr_html_renderer = ""
            if os.environ.get("AE_DISABLE_HTML_SLIDES", "").strip() == "1":
                _rr_html_reason = "AE_DISABLE_HTML_SLIDES=1"
            elif not _html_on:
                _rr_html_reason = "opted-out (editable/native requested)"
            else:
                _rr_html_reason = "bridge-unreachable & no-local-chrome"

        def _embed_fullbleed(slide, png_abs):
            """PNG를 슬라이드 전체(13.333x7.5in)에 배경으로 깐다(back-most)."""
            try:
                pic = slide.shapes.add_picture(png_abs, Inches(0), Inches(0),
                                               width=Inches(13.333), height=Inches(7.5))
                spTree = pic._element.getparent()
                spTree.remove(pic._element)
                spTree.insert(2, pic._element)
                return True
            except Exception as _ee:
                print(f"[generate_pptx] full-bleed embed 실패: {str(_ee)[:160]}")
                return False

        def _next_slide(layout):
            """다음 대상 슬라이드를 반환. 템플릿이면 도너 재사용/복제, 아니면 레이아웃에서 생성."""
            if _use_designed:
                idx = _donor_consumed["n"]
                if idx < len(_donors):
                    sl = _donors[idx]
                else:
                    sl = _clone_slide(prs, _content_donor) or _safe_add_slide(prs, layout)
                _donor_consumed["n"] = idx + 1
                _strip_slide_to_background(sl)
                return sl
            return _safe_add_slide(prs, layout)

        # Cover slide
        # 요구사항 6.2 — 표지는 "title" 레이아웃으로 해석. 무템플릿이면 _resolve_layout이
        # LAYOUT_MAP["title"]=0 → slide_layouts[0]을 그대로 반환해 기존 동작과 동일하다(5.2).
        cover_layout = _resolve_layout(prs, "title", used_template)
        cover = _next_slide(cover_layout)
        # 요구사항 6.8 — 표지 제목/날짜는 placeholder text_frame에 편집 가능한 텍스트로
        # 채운다(이미지 래스터화 금지). 템플릿 적용 여부와 무관하게 PowerPoint에서 편집 가능.
        # 요구사항 9.5 — title placeholder 부재 등 템플릿 편차는 _safe_set_title이 격리한다.
        _safe_set_title(cover, title)
        from datetime import datetime as _dt
        _cover_sub = _dt.now().strftime("%Y-%m-%d")
        _sub_set = False
        if len(cover.placeholders) > 1:
            # 요구사항 9.5 — 부제/날짜 placeholder 설정은 템플릿 단계로 간주해 격리한다.
            try:
                cover.placeholders[1].text = _cover_sub
                _sub_set = True
            except Exception as e:
                print(f"[generate_pptx] cover date placeholder 설정 실패: {str(e)[:200]}")
        if not _sub_set:
            # 부제 placeholder 부재 — 제목 아래 명시적 텍스트박스로 날짜를 넣어
            # 표지가 제목만 남지 않게 한다(이슈1).
            try:
                from pptx.util import Pt as _Pt
                from pptx.dml.color import RGBColor as _RGB
                _sb = cover.shapes.add_textbox(Inches(0.7), Inches(1.45), Inches(11.9), Inches(0.6))
                _sp = _sb.text_frame.paragraphs[0]
                _sr = _sp.add_run()
                _sr.text = _cover_sub
                _sr.font.size = _Pt(18)
                _sr.font.color.rgb = _RGB(0x6A, 0x6A, 0x6A)
                try:
                    _sr.font.name = "Malgun Gothic"
                except Exception:
                    pass
            except Exception as e:
                print(f"[generate_pptx] cover subtitle textbox 실패: {str(e)[:160]}")
        # Optional cover background — caller passes via tool_input["coverBackground"].
        # 요구사항 6.5 — full-bleed 표지 배경(1920×1080)이 있으면 (0,0)에 13.333×7.5in로
        #   add_picture 후 spTree.insert(2, ...)로 마스터 배경 위·placeholder 아래에 배치해
        #   템플릿 마스터 배경을 가린다. 0=nvGrpSpPr, 1=grpSpPr 이므로 인덱스 2가 첫 shape
        #   자리(back-most). add_slide가 먼저 복제한 placeholder들은 3+로 밀려 그림 위에 남는다.
        # 요구사항 6.6 — coverBackground가 없거나 파일이 없으면 그림을 추가하지 않으므로
        #   템플릿 마스터 배경이 그대로 보인다.
        cover_bg = tool_input.get("coverBackground", "") if isinstance(tool_input, dict) else ""
        _cover_bg_embedded = False
        if isinstance(cover_bg, str) and cover_bg.strip():
            cand = cover_bg if os.path.isabs(cover_bg) else os.path.join(_local_root, cover_bg.strip())
            # 경로 어긋남 방지 — _local_root join이 빗나가면 다중 루트 탐색으로 재해석
            # (생성된 Vertex/hero 이미지가 경로 불일치로 고아되는 반복 문제 해결).
            if not os.path.isfile(cand):
                try:
                    _alt = _resolve_relative_for_verify(cover_bg.strip(), project_path)
                    if _alt and os.path.isfile(_alt):
                        cand = _alt
                    else:
                        _bn = os.path.basename(cover_bg.strip())
                        _g = os.path.join(_local_root, ".generated", _bn)
                        if os.path.isfile(_g):
                            cand = _g
                except Exception:
                    pass
            if os.path.isfile(cand):
                try:
                    pic = cover.shapes.add_picture(cand, Inches(0), Inches(0),
                                                    width=Inches(13.333), height=Inches(7.5))
                    spTree = pic._element.getparent()
                    spTree.remove(pic._element)
                    spTree.insert(2, pic._element)  # 0=nvGrpSpPr, 1=grpSpPr, 2+=shapes
                    _cover_bg_embedded = True
                except Exception as e:
                    print(f"[generate_pptx] cover bg embed failed: {e}")
        # 배경 이미지가 없으면 편집 가능한 네이티브 디자인 표지 적용(템플릿 표지 거의 빈
        # 문제 해결). AE_DISABLE_NATIVE_COVER=1로 비활성(기존 placeholder 표지 폴백).
        _cover_native_applied = False  # task20(C): 네이티브 표지 적용 여부
        if os.environ.get("AE_DISABLE_NATIVE_COVER", "") != "1" and not _cover_bg_embedded:  # B방향(수정1): 통짜 표지 배경 채택 시 네이티브 표지 억제(공존 금지)
            try:
                from native_diagram_pptx import build_native_cover as _bnc
            except Exception:
                try:
                    from ai_engine.native_diagram_pptx import build_native_cover as _bnc
                except Exception:
                    _bnc = None
            if _bnc is not None:
                _cv_sub = ""
                _cv_eb = ""
                if isinstance(tool_input, dict):
                    _cv_sub = str(tool_input.get("subtitle", "") or "")
                    _cv_eb = str(tool_input.get("eyebrow", "") or "")
                if not _cv_sub and slides_data:
                    try:
                        _fs = slides_data[0]
                        if isinstance(_fs, dict):
                            _cv_sub = str(_fs.get("subtitle", "") or _fs.get("title", "") or "")
                    except Exception:
                        pass
                # KPI 추출(선택) — kpi로 분류되는 슬라이드가 있으면 표지 요약 카드로.
                _cv_kpis = None
                try:
                    try:
                        from native_diagram_pptx import _parse_kpis as _pk
                    except Exception:
                        from ai_engine.native_diagram_pptx import _parse_kpis as _pk
                    for _sd in (slides_data or [])[:6]:
                        if not isinstance(_sd, dict):
                            continue
                        _h = str(_sd.get("title", "") or "")
                        _b = "\n".join(str(x) for x in (_sd.get("bullets", []) or []))
                        _k, _c = _classify_section_diagram(_h, _b, title)
                        if _k == "kpi":
                            _kp = _pk(_c or _b)
                            if _kp:
                                _cv_kpis = _kp
                                break
                except Exception:
                    _cv_kpis = None
                try:
                    from datetime import datetime as _dt2
                    _ok_cv = _bnc(cover, prs, title, subtitle=_cv_sub, eyebrow=_cv_eb,
                                  date_str=_dt2.now().strftime("%Y-%m-%d"),
                                  kpis=_cv_kpis, palette=_tpl_palette_for_native,
                                  remove_placeholders=True,
                                  over_image=_cover_bg_embedded)
                    if _ok_cv:
                        _cover_native_applied = True  # task20(C)
                        print(f"[generate_pptx] 표지 네이티브 디자인 적용(kpis={'Y' if _cv_kpis else 'N'})")
                        # 수정 2(게이트 보강): 표지 density/style chrome. build_native_cover
                        # 가 못 채우는 icon_badge/step_grid/accent 색을 design_tokens
                        # (산출물_검증기와 동일 출처) 기준으로 가산 방출해 표지도 (e)밀도/
                        # (h)스타일 audit 합격. 텍스트 없는 장식 → 겹침 검사 제외·콘텐츠
                        # 뒤 z-order. KPI 카드가 이미 있으면 step_grid 중복을 피한다.
                        try:
                            try:
                                from ai_engine.slide_templates import design_tokens_for_profile as _dtfp_cc
                            except ImportError:
                                from slide_templates import design_tokens_for_profile as _dtfp_cc
                            try:
                                from ai_engine.native_layout_renderer import emit_cover_density_chrome as _eccd
                            except ImportError:
                                from native_layout_renderer import emit_cover_density_chrome as _eccd
                            _eccd(cover, _dtfp_cc(style_profile), add_step_grid=not bool(_cv_kpis))
                        except Exception as _ecce:
                            print(f"[generate_pptx] 표지 chrome 보강 실패(무시): {str(_ecce)[:160]}")
                except Exception as e:
                    print(f"[generate_pptx] native cover 실패: {str(e)[:200]}")
        # 빈 placeholder(프롬프트/점선) 제거 — 표지를 깔끔하게.
        # B방향(수정1): 통짜 표지 배경(_cover_bg_embedded)이 채택되면 그 위의
        # 편집 제목/부제/날짜/네이티브 표지 텍스트·도형을 모두 제거해 '통짜 하나'만 남긴다
        # (겹츨0·중복0). 배경(HTML 표지/Vertex)이 제목을 포함한다고 신뢰.
        if _cover_bg_embedded:
            _strip_text_over_fullbleed(cover)
        _remove_empty_placeholders(cover)
        _remove_empty_text_shapes(cover)
        _clamp_shapes_into_bounds(cover)  # B방향(수읅5): 표지 도형 경계 후처리 클램프

        # 표지도 Genspark급 HTML 디자인으로 (활성 시). 실패하면 위 placeholder 제목 유지.
        if _html_enabled:
            try:
                from datetime import datetime as _dtc
                try:
                    from ai_engine.slide_templates import render_layout as _rl_cov, design_tokens_for_profile as _dtfp_cov
                except ImportError:
                    from slide_templates import render_layout as _rl_cov, design_tokens_for_profile as _dtfp_cov
                _cov_data = {"title": title or "", "subtitle": _dtc.now().strftime("%Y-%m-%d"),
                             "eyebrow": "PRESENTATION"}
                if style_profile:
                    try:
                        _cov_data["design"] = _dtfp_cov(style_profile)
                    except Exception:
                        pass
                _cov_html = _rl_cov("cover", _cov_data)
                if _cov_html:
                    import time as _tc_cov
                    _cov_abs = os.path.join(_local_root, ".generated", f"cover-{int(_tc_cov.time()*1000)}.png")
                    os.makedirs(os.path.dirname(_cov_abs), exist_ok=True)
                    _cres = await _render_html_slide_to_png(_cov_html, _cov_abs, 1920, 1080, 30)
                    # D1(슬라이드당 풀블리드 1회): coverBackground 가 이미 풀블리드로
                    # 임베드됐으면(_cover_bg_embedded) HTML 표지 풀블리드를 스킵한다.
                    # 손실-0: 스킵된 HTML 표지 PNG(_cov_abs)는 폐기하지 않고 디스크에 보존.
                    if (_cres.get("ok") and os.path.isfile(_cov_abs)
                            and not _cover_bg_embedded
                            and not _cover_native_applied
                            and os.environ.get("AE_COVER_HTML_FULLBLEED", "0") == "1"):  # task20 수정B: 콘텐츠 구운 HTML 표지 풀블리드 미채택(기본 OFF, 킬스위치=1로 복원) / task20(C): 네이티브 표지 우선
                        _embed_fullbleed(cover, _cov_abs)
                        _rr_cover_path = "html-fullbleed"
                        try:
                            _remove_empty_text_shapes(cover)
                        except Exception:
                            pass
            except Exception as _che:
                print(f"[generate_pptx] 표지 HTML 실패(네이티브 유지): {str(_che)[:160]}")

        # === Vertex(Nano Banana Pro) 이미지 병렬 일괄 생성 (블로킹/끊김 방지) ===
        # 슬라이드마다 순차 await(최대 60s)하면 긴 덱에서 수 분간 멈춘 것처럼 보여
        # '무한루프/끊김'으로 인식된다. 진입 전에 병렬로 한 번에 생성하고 루프에서는
        # 결과만 조회한다. HTML 디자인 슬라이드가 활성이면(_html_enabled) 그쪽이
        # 풀블리드를 담당하므로 Vertex 일괄 생성은 건너뛴다.
        # === 편집 가능 네이티브 다이어그램 우선 (젤스파크 방식) ===
        # 모든 도형이 PowerPoint에서 편집 가능해야 한다. 섹션을 LLM으로 구조화해
        # sd['nativeDiagram']에 기록하면 네이티브 도형으로 렌더되고, 아래 Vertex는 스킵된다.
        if (not _html_enabled) and os.environ.get("AE_PREFER_EDITABLE_DIAGRAM", "1") != "0":
            try:
                _struct_gw = _get_gw(aws_profile, bedrock_user)
                _struct_model = _specialized_model_for_task(
                    "file_generation", "", aws_profile=aws_profile, bedrock_user=bedrock_user)
            except Exception:
                _struct_gw = None
                _struct_model = ""
            if _struct_gw and _struct_model:
                _doc_ctx_s = (title or "")[:300]

                async def _struct_one(_idx, _sd):
                    if not isinstance(_sd, dict):
                        return _idx, None
                    if _sd.get("slideBackground") or _sd.get("imageFile") or _sd.get("nativeDiagram"):
                        return _idx, None
                    _h = str(_sd.get("title", "") or "").strip()
                    _b = "\n".join(str(x) for x in (_sd.get("bullets", []) or []))
                    if not (_h or _b):
                        return _idx, None
                    try:
                        spec = await _llm_structure_native_diagram(_struct_gw, _struct_model, _h, _b, _doc_ctx_s)
                    except Exception as _e:
                        print(f"[generate_pptx] 구조화 예외 slide {_idx + 2}: {str(_e)[:140]}")
                        return _idx, None
                    return _idx, (spec or None)

                _selig = list(enumerate(slides_data))[:12]
                if _selig:
                    print(f"[generate_pptx] 편집가능 다이어그램 구조화 대상 {len(_selig)}장")
                    try:
                        _sres = await asyncio.gather(
                            *(_struct_one(_i, _s) for _i, _s in _selig),
                            return_exceptions=True)
                        _sok = 0
                        for _r in _sres:
                            if isinstance(_r, Exception):
                                continue
                            _ri, _spec = _r
                            if _spec and isinstance(slides_data[_ri], dict):
                                slides_data[_ri]["nativeDiagram"] = _spec
                                _sok += 1
                        print(f"[generate_pptx] 편집가능 다이어그램 구조화 성공 {_sok}장")
                    except Exception as _se:
                        print(f"[generate_pptx] 구조화 일괄 예외: {str(_se)[:160]}")

            # 결정론적 폴백 — LLM 구조화가 비거나 게이트웨이 미가용이어도 본문을
            # '휘한 불릿' 대신 편집 가능한 네이티브 카드 그리드로 채운다(요구: 본문도 고퀸).
            if (not _html_enabled) and os.environ.get("AE_PREFER_EDITABLE_DIAGRAM", "1") != "0":
                import re as _re_cf
                def _norm_cards(_bullets):
                    """불릿/문단을 카드 줄 리스트로 정규화. 마크다운 제거, 긴 문단은 분할."""
                    _items = []
                    for _b in (_bullets or []):
                        _t = str(_b).strip()
                        if not _t:
                            continue
                        _t = _re_cf.sub(r"`([^`]+)`", r"\1", _t)
                        _t = _t.replace("**", "").replace("__", "")
                        _t = _re_cf.sub(r"^\s*(?:[-*•]\s+|#{1,6}\s+)", "", _t)
                        if _t:
                            _items.append(_t)
                    # 한 문단뽐이면 문장/구분자로 분할해 카드 그리드로
                    if len(_items) == 1 and len(_items[0]) > 40:
                        _segs = _re_cf.split(r"(?:\n|·|•|;|\. |。)", _items[0])
                        _segs = [s.strip(" .") for s in _segs if len(s.strip()) > 3]
                        if len(_segs) >= 2:
                            _items = _segs
                    return _items[:6]
                try:
                    for _ci, _cs in list(enumerate(slides_data))[:12]:
                        if not isinstance(_cs, dict):
                            continue
                        if _cs.get("slideBackground") or _cs.get("imageFile") or _cs.get("nativeDiagram"):
                            continue
                        _cbl = _norm_cards(_cs.get("bullets", []))
                        if len(_cbl) >= 4:
                            _cs["nativeDiagram"] = {"type": "twocol", "content": "\n".join(_cbl)}
                        elif len(_cbl) >= 2:
                            _cs["nativeDiagram"] = {"type": "cards", "content": "\n".join(_cbl)}
                except Exception as _cfe:
                    print(f"[generate_pptx] 카드 폴백 예외: {str(_cfe)[:120]}")

        _vertex_pre = {}
        # Vertex 사전생성 게이트 — HTML과 공존(task 3.2). 과거 `not _html_enabled` 상호배타
        # 게이트는 HTML 고품질 경로가 켜지면 Vertex 이미지를 아예 생성하지 못하게 만들어
        # "HTML 품질"과 "Vertex 이미지"가 구조적으로 공존 불가였다(Bug Condition gateSuppressed).
        # 이제 _html_enabled 여부와 무관하게 실행한다. 효과적 게이트는
        # `AE_PREFER_VERTEX_IMAGE != 0 AND vertexClient.enabled`(enabled 판정은 아래 _vc_pptx).
        # 역할이 structural인 슬라이드만 _gen_vertex_slide 내부에서 스킵(네이티브 도형 우선, Req 3.1).
        if os.environ.get("AE_PREFER_VERTEX_IMAGE", "1") != "0":
            try:
                try:
                    from ai_engine.vertex_image_module import get_vertex_image_client as _vget_pptx
                except ImportError:
                    from vertex_image_module import get_vertex_image_client as _vget_pptx
                _vc_pptx = _vget_pptx(aws_profile=aws_profile or "")
            except Exception:
                _vc_pptx = None
            if _vc_pptx and getattr(_vc_pptx, "enabled", False):
                _rr_vertex_enabled = True
                import time as _vt, hashlib as _vh, base64 as _vb
                _vgd = os.path.join(_local_root, ".generated")
                os.makedirs(_vgd, exist_ok=True)

                async def _gen_vertex_slide(_idx, _sd):
                    if not isinstance(_sd, dict):
                        return _idx, ""
                    if _sd.get("slideBackground") or _sd.get("imageFile") or _sd.get("nativeDiagram"):
                        return _idx, ""
                    # 역할 기반 스킵 통일(task 3.2) — structural(흐름/트리/아키텍처) 슬라이드만
                    # 네이티브 편집 도형 우선으로 Vertex 래스터 생성을 건너뛴다(Req 3.1 보존).
                    # cover/content/visual 슬라이드는 Vertex 이미지 생성 대상으로 남는다.
                    # (과거: _classify_section_diagram이 낸 임의 kind(kpi/cards/block 등)까지
                    #  스킵해 고밀도 콘텐츠/비주얼 슬라이드의 이미지가 부당하게 누락됐다.)
                    _role = "content"
                    try:
                        _role = _classify_slide_role(_sd, False, title)
                        if _role == "structural":
                            return _idx, ""
                    except Exception:
                        _role = "content"
                    _vtitle = str(_sd.get("title", "") or "")[:120]
                    _vbul = [str(b).strip() for b in (_sd.get("bullets", []) or []) if str(b).strip()]
                    _vctx = " ".join(_vbul)[:500]
                    if not (_vtitle or _vctx):
                        return _idx, ""
                    # 역할 기반 프롬프트(task 3.2) — structural은 위에서 스킵되었으므로 여기
                    # 도달하는 슬라이드는 visual(히어로/사진형) 또는 content(고밀도 콘텐츠)다.
                    #   visual  → 사진·일러스트형 히어로 비주얼(텍스트 없는 분위기 이미지)
                    #   content → HTML 고밀도 레이아웃의 이미지 슬롯에 들어갈 보조 비주얼
                    # 본문 텍스트·구조·라벨은 HTML/네이티브가 담당하므로 이미지에는 글자를 넣지
                    # 않는다(Nano Banana Pro의 텍스트 렌더 강점은 표지/네이티브 경로에서 활용).
                    _topic = (_vtitle + " " + _vctx).strip()
                    # 기본(레거시) negative — 하이브리드 OFF 또는 content 경로에서 사용(바이트 동일 보존).
                    _vnegative = ("watermark, fake logo, brand name, emoji, distorted text, "
                                  "unreadable artifacts, childish clipart")
                    if _hybrid_on and _role in ("cover", "section", "visual"):
                        # 하이브리드 풀블리드 대상(cover/section/visual) — 프롬프트 빌더(task 4.1/7.3)로
                        # 역할별 프롬프트와 no-text negative를 결정론적으로 조립(Req 3.1/3.2/3.5/3.7).
                        _vprompt, _vnegative = _build_fullbleed_vertex_prompt(
                            _role, _vtitle, _vbul, style_profile)
                    elif _role == "visual":
                        _vprompt = (
                            f'A premium photographic or polished editorial illustration serving as the hero '
                            f'visual for a corporate presentation slide. Theme: "{_vtitle}". Context: {_topic}. '
                            f'Style: cinematic professional photography or clean modern flat illustration, '
                            f'natural soft lighting, shallow depth of field, refined corporate aesthetic, '
                            f'generous negative space for overlaid text, balanced 16:9 composition. '
                            f'NO text, NO words, NO letters, NO captions, NO charts, NO diagrams, NO emoji, '
                            f'NO watermark, NO fake logos.'
                        )
                    else:
                        _vprompt = (
                            f'A clean auxiliary supporting visual to accompany a text-rich content slide '
                            f'titled "{_vtitle}". Concept/context: {_topic}. '
                            f'Style: modern flat corporate illustration or conceptual imagery that complements '
                            f'a dense layout, professional blue and navy palette with one accent color, ample '
                            f'whitespace, balanced composition suited to an image slot beside text, 16:9. '
                            f'NO embedded text, NO words, NO letters, NO charts, NO emoji, NO watermark, '
                            f'NO fake logos.'
                        )
                    try:
                        _vres = await _vc_pptx.generate(
                            prompt=_vprompt, model_class="image_generation_high_quality",
                            aspect_ratio="16:9",
                            negative_prompt=_vnegative,
                            timeout=60,
                        )
                    except Exception as _e:
                        print(f"[generate_pptx] Vertex 예외 slide {_idx + 2}(무시): {str(_e)[:140]}")
                        return _idx, ""
                    _imgs = (_vres.get("images") or []) if isinstance(_vres, dict) else []
                    if not _imgs:
                        if isinstance(_vres, dict) and _vres.get("error"):
                            print(f"[generate_pptx] Vertex 실패 slide {_idx + 2}: {_vres.get('error')} -> 네이티브 폴백")
                        return _idx, ""
                    try:
                        _fn = f"vertex-slide-{int(_vt.time()*1000)}-{_vh.md5((_vprompt+str(_idx)).encode()).hexdigest()[:6]}-{_idx}.png"
                        with open(os.path.join(_vgd, _fn), "wb") as _f:
                            _f.write(_vb.b64decode(_imgs[0]))
                        return _idx, f".generated/{_fn}"
                    except Exception as _e2:
                        print(f"[generate_pptx] Vertex 저장 실패 slide {_idx + 2}: {str(_e2)[:140]}")
                        return _idx, ""

                _elig = list(enumerate(slides_data))[:12]
                if _elig:
                    print(f"[generate_pptx] Vertex 이미지 병렬 생성 — 대상 {len(_elig)}장")
                    # 동시 호출을 3개로 제한 — Gemini 이미지 분당 쿼터로 인한
                    # http-429 다발을 방지(초과분은 네이티브 폴백됨). 재시도는 vertex 모듈이 처리.
                    _vsem = asyncio.Semaphore(3)
                    async def _bounded_vertex(_bi, _bs):
                        async with _vsem:
                            return await _gen_vertex_slide(_bi, _bs)
                    try:
                        _vresults = await asyncio.gather(
                            *(_bounded_vertex(_i, _s) for _i, _s in _elig),
                            return_exceptions=True)
                        for _r in _vresults:
                            if isinstance(_r, Exception):
                                continue
                            _ri, _rel = _r
                            if _rel:
                                _vertex_pre[_ri] = _rel
                    except Exception as _ge:
                        print(f"[generate_pptx] Vertex 일괄 생성 예외(무시): {str(_ge)[:160]}")
                    print(f"[generate_pptx] Vertex 이미지 생성 완료 — 성공 {len(_vertex_pre)}장")

        # === renderReport: Vertex 비활성 사유 판정 (관측만) ===
        if not _rr_vertex_enabled and not _rr_vertex_reason:
            if os.environ.get("AE_PREFER_VERTEX_IMAGE", "1") == "0":
                _rr_vertex_reason = "AE_PREFER_VERTEX_IMAGE=0"
            elif os.environ.get("AE_DISABLE_VERTEX_IMAGE", "").strip() == "1":
                _rr_vertex_reason = "AE_DISABLE_VERTEX_IMAGE=1"
            else:
                _rr_vertex_reason = "no key resolved"

        # === 목차(TOC) 자동 생성 — 슬라이드 덱의 기본. 편집 가능 네이티브 카드로. ===
        # 콘텐츠 제목이 3개 이상이면 표지 다음에 목차를 넣는다. AE_PPTX_TOC=0으로 비활성.
        if (not _html_enabled) and os.environ.get("AE_PPTX_TOC", "1") != "0":
            _toc_titles = []
            for _sd in slides_data:
                if isinstance(_sd, dict):
                    _t = str(_sd.get("title", "") or "").strip()
                    if _t:
                        _toc_titles.append(_t)
            if len(_toc_titles) >= 3:
                try:
                    _toc_layout = _resolve_layout(prs, "content", used_template)
                    _toc_slide = _next_slide(_toc_layout)
                    _safe_set_title(_toc_slide, "목차")
                    try:
                        if len(_toc_slide.placeholders) > 1:
                            _toc_slide.placeholders[1].text_frame.clear()
                    except Exception:
                        pass
                    try:
                        from native_diagram_pptx import build_native_diagram as _bnd_toc
                    except Exception:
                        try:
                            from ai_engine.native_diagram_pptx import build_native_diagram as _bnd_toc
                        except Exception:
                            _bnd_toc = None
                    if _bnd_toc is not None:
                        _bnd_toc(_toc_slide, "block", "\n".join(_toc_titles[:8]),
                                 region=(0.7, 1.7, 11.9, 5.2),
                                 palette=_tpl_palette_for_native, title="")
                    try:
                        _remove_empty_text_shapes(_toc_slide)
                    except Exception:
                        pass
                    print(f"[generate_pptx] 목차 슬라이드 생성 — {len(_toc_titles)}개 항목")
                except Exception as _te:
                    print(f"[generate_pptx] 목차 생성 실패(무시): {str(_te)[:160]}")

        # task 8.5 — 하이브리드 content 편집 경로로 렌더된 슬라이드의 id 집합. 덱 레벨
        # chrome 패스(_apply_universal_slide_design)에서 이 슬라이드들을 건너뛰어 native
        # 렌더러 자족 chrome 과의 이중 방출(제목/푸터 겹침)을 방지한다(단일 chrome 소스).
        _hybrid_content_slide_ids = set()
        for i, sd in enumerate(slides_data):
            if not isinstance(sd, dict):
                sd = {"title": str(sd)}
            layout_name = sd.get("layout", "content")
            # 요구사항 6.2/6.3/6.4 — 템플릿이면 동적 매핑(이름→콘텐츠 레이아웃→index 0),
            # 무템플릿이면 기존 LAYOUT_MAP 동작 보존(5.2).
            layout = _resolve_layout(prs, layout_name, used_template)
            s = _next_slide(layout)
            # 요구사항 6.8 — 제목/본문은 placeholder text_frame에 편집 가능한 텍스트로 채운다
            # (이미지 래스터화 금지). 템플릿 적용 여부와 무관하게 PowerPoint에서 편집 가능 유지.
            # 요구사항 9.5 — title placeholder 부재 등 템플릿 편차는 _safe_set_title이 격리한다.
            # task 3.1/3.2 — 콘텐츠가 구워진 풀블리드 HTML 배경(slideBackground)이 채택되면
            # 그 PNG가 제목/본문을 이미 포함하므로(baked-in) 네이티브 제목을 중복 방출하지
            # 않는다(제목 1회). HTML 베이크 후보가 아니면 기존과 동일하게 즉시 제목을 설정해
            # 비결함 경로를 그대로 보존한다(additive). 베이크 실패 시 게이트 직후 폴백 설정.
            _html_bake_eligible = bool(
                _html_enabled and not sd.get("slideBackground")
                and not sd.get("imageFile") and not sd.get("nativeDiagram"))
            if not _html_bake_eligible:
                _safe_set_title(s, sd.get("title", f"Slide {i + 2}"))

            bullets = sd.get("bullets", [])
            # task 3.5 — Vertex 이미지가 HTML 레이아웃에 합성되었는지 추적(손실 0 판정).
            _pre_composited_into_html = False
            # renderReport per-slide 관측 상태 (caller 명시값은 mutate 전에 포착)
            _rr_caller_img = bool(sd.get("imageFile"))
            _rr_caller_bg = bool(sd.get("slideBackground"))
            _rr_html_bg_set = False
            _rr_plan_slot = ""
            _rr_role = ""
            # task 7.2 — 하이브리드 주 렌더러 선택(플랜 결정) 관측 상태. 기본값은
            # flag OFF/비대상 슬라이드에서 그대로 유지된다(가산적). 실제 렌더 스왑은 7.3/7.4.
            _hybrid_plan = None
            _rr_hybrid_primary = ""
            _rr_hybrid_slot = ""
            _rr_hybrid_editable = None
            # task 7.4 — 하이브리드 content 편집 경로로 라우팅되었는지 표시(게이팅·per-slide).
            # True면 (a) HTML 풀블리드 바이크를 우회했고 (b) 아래 Vertex 임베드 블록이
            # 히어로를 풀블리드로 재임베드하지 않도록 스킵해 content 풀블리드 PICTURE 0을 보존한다.
            _hybrid_content_routed = False
            # Genspark급 HTML 디자인 슬라이드 — 활성 시 이 슬라이드를 풀블리드 배경으로 렌더.
            # 성공하면 slideBackground가 설정되어 네이티브 다이어그램 추론이 자동 skip된다.
            # === task 9.1 — 네이티브 라우팅 게이트 (design §Components §2, Property 3/6) ===
            # 알려진_레이아웃 + caller 비명시 + 콘텐츠 텍스트면 콘텐츠 베이크(HTML→PNG→
            # slideBackground) 대신 편집가능 네이티브 도형으로 렌더한다(Req 1.5/5.1).
            # 가산적·보수적: AE_NATIVE_LAYOUT_RENDER!=1 이면 즉시 비활성(기존 베이크 보존),
            # 게이트 False(비-네이티브 레이아웃/명시 경로)면 기존 흐름 그대로(no-op, Req 9).
            # 네이티브 렌더/폴백 실패 시 예외를 흙수하고 베이크 경로로 폴백(콘텐츠 손실 0).
            # === task20 수정 A — 본문 콘텐츠 네이티브 무조건 라우팅(베이크 원천 차단) ===
            # AE_NATIVE_LAYOUT_RENDER!=0(기본 활성) + caller 비명시 + 콘텐츠 존재면
            # 게이트웨이 픽 결과와 무관하게 반드시 편집가능 네이티브로 렌더하고
            # _native_routed=True 를 보장한다. 게이트웨이 미가용/타임아웃/빈 결과여도
            # 합성 레이아웃(feature_grid/two_column/section_divider)으로 렌더 → 아래
            # HTML->PNG 베이크 블록은 not _native_routed 가드로 자동 skip(통짜 금지, Req 1.5).
            _native_routed = False
            _nl_has_content = bool(str(sd.get("title", "") or "").strip()) or any(
                str(_b).strip() for _b in (bullets or []))
            if (_html_enabled
                    and os.environ.get("AE_NATIVE_LAYOUT_RENDER", "0") == "1"  # B방향: 기본 OFF(통짜 우선), 명시 옵트인(=1)만 네이티브 라우팅
                    and _nl_has_content
                    and not sd.get("slideBackground") and not sd.get("imageFile")
                    and not sd.get("nativeDiagram")):
                _nl_title = str(sd.get("title", "") or "")
                _nl_body = "\n".join(str(_b) for _b in bullets)
                _nl_bl = [str(_b).strip() for _b in (bullets or []) if str(_b).strip()]
                # 1) 게이트웨이 픽(있으면). 미가용/타임아웃/예외/빈 결과는 흡수 —
                #    실패해도 아래 합성 레이아웃으로 계속(네이티브 라우팅 강제).
                _nl_pick_layout = ""
                _nl_pick_data = {}
                try:
                    _nl_role = _classify_slide_role(sd, False, title)
                except Exception:
                    _nl_role = "content"
                try:
                    _nl_bn = len([_l for _l in _nl_body.splitlines() if _l.strip()])
                    _nl_pick = await _llm_pick_slide_layout(
                        _html_gw, _html_model, str(sd.get("title", "")), _nl_body,
                        title, role=_nl_role, bullet_count=_nl_bn)
                    _nl_pick_layout = (_nl_pick or {}).get("layout") or ""
                    _nl_pick_data = (_nl_pick or {}).get("data") or {}
                except Exception as _pke:
                    print(f"[generate_pptx] 레이아웃 픽 실패(합성 레이아웃 사용) slide {i + 2}: {str(_pke)[:140]}")
                    _nl_pick_layout = ""
                    _nl_pick_data = {}
                # 2) 렌더러/토큰 import (실패해도 폴백 계속).
                try:
                    from ai_engine.native_layout_renderer import (
                        map_to_native_layout as _m2nl,
                        render_native_layout as _rnl,
                        render_native_fallback as _rnf)
                except ImportError:
                    from native_layout_renderer import (
                        map_to_native_layout as _m2nl,
                        render_native_layout as _rnl,
                        render_native_fallback as _rnf)
                try:
                    from ai_engine.slide_templates import design_tokens_for_profile as _dtfp_nl
                except ImportError:
                    from slide_templates import design_tokens_for_profile as _dtfp_nl
                try:
                    _nl_tokens = _dtfp_nl(style_profile)
                except Exception:
                    _nl_tokens = {}
                _nl_used = ""
                try:
                    # 3) 렌더 후보 사다리 — 게이트웨이 픽 → 합성(제목+불릿) → 최소.
                    #    각 후보를 render_native_layout 로 시도(ok=True 면 chrome 포함
                    #    리치 렌더 → 밀도/스타일 audit 통과). 어떤 픽/빈결과든 반드시
                    #    편집가능 네이티브로 귀결(통짜 이미지 금지, Req 1.5/5.1).
                    _nl_cands = []
                    _nl_mapped = _m2nl(_nl_pick_layout) if _nl_pick_layout else ""
                    if _nl_mapped and isinstance(_nl_pick_data, dict) and _nl_pick_data:
                        _nl_cands.append((_nl_mapped, _nl_pick_data))
                    if _nl_bl:
                        _nl_cands.append(("feature_grid",
                                          {"title": _nl_title,
                                           "features": [{"title": _b, "description": ""} for _b in _nl_bl]}))
                        _nl_half = (len(_nl_bl) + 1) // 2
                        _nl_cands.append(("two_column",
                                          {"title": _nl_title,
                                           "left_content": _nl_bl[:_nl_half] or [_nl_title],
                                           "right_content": _nl_bl[_nl_half:] or [_nl_title]}))
                    _nl_cands.append(("section_divider", {"title": _nl_title or f"Slide {i + 2}"}))
                    for _nl_cl, _nl_cd in _nl_cands:
                        try:
                            _nl_res = _rnl(s, prs, _nl_cl, _nl_cd, _nl_tokens,
                                           aws_profile=aws_profile, credentials=None)
                        except Exception as _rle:
                            print(f"[generate_pptx] 네이티브 렌더 예외(다음 후보) slide {i + 2} layout={_nl_cl}: {str(_rle)[:140]}")
                            _nl_res = None
                        if _nl_res is not None and getattr(_nl_res, "ok", False):
                            _native_routed = True
                            _nl_used = _nl_cl
                            break
                    # 4) 리치 후보 전부 실패 → 편집가능 텍스트 폴백(통짜 금지, Req 1.4).
                    if not _native_routed:
                        try:
                            _nl_fb = _rnf(s, {"title": _nl_title, "bullets": list(bullets or [])}, _nl_tokens)
                            if _nl_fb is not None and getattr(_nl_fb, "ok", False):
                                _native_routed = True
                                _nl_used = "fallback"
                        except Exception as _fbe:
                            print(f"[generate_pptx] 네이티브 폴백 실패 slide {i + 2}: {str(_fbe)[:140]}")
                    # 5) 최후 보증 — 콘텐츠 슬라이드는 절대 베이크로 안 감. 폴백조차
                    #    실패했더라도 편집가능 제목 텍스트박스만이라도 두고 skip 한다.
                    if not _native_routed:
                        try:
                            try:
                                from ai_engine.native_layout_renderer import emit_title as _et_fb, _TITLE_REGION as _tr_fb
                            except ImportError:
                                from native_layout_renderer import emit_title as _et_fb, _TITLE_REGION as _tr_fb
                            _et_fb(s, _nl_title or f"Slide {i + 2}", _nl_tokens, _tr_fb)
                            _nl_used = "min-title"
                        except Exception:
                            pass
                        _native_routed = True
                except Exception as _nle:
                    # 라우팅 자체 예외 — 그래도 콘텐츠 슬라이드는 베이크로 보내지 않는다.
                    print(f"[generate_pptx] 네이티브 라우팅 예외(베이크 미전환) slide {i + 2}: {str(_nle)[:160]}")
                    _native_routed = True
                if _native_routed:
                    print(f"[generate_pptx] 슬라이드 {i + 2} → 네이티브 레이아웃 렌더(layout={_nl_used or 'native'})")

            # === task 7.4: 하이브리드 content 편집 경로 배선 (게이팅·가산적) ===
            # `_hybrid_on == True`이고 caller가 imageFile/slideBackground/nativeDiagram을
            # 지정하지 않았으며 아직 네이티브 라우팅되지 않은 슬라이드에 한해, 역할을 예외
            # 안전하게 확정(예외 시 "content")한다. role=="content"(하이브리드 플랜 primary=
            # NATIVE_EDITABLE)이면 _generate_html_slide_for_section 풀블리드 PNG 바이크를
            # 우회하고 _render_content_editable(편집 가능 네이티브 + 바운디드 히어로 합성/보존)로
            # 라우팅해 _native_routed=True로 확정한다(R2.1/R2.2/R2.3/R2.4). Vertex 히어로는
            # _vertex_pre[i]가 있으면 바운디드 슬롯에 합성, 없으면 네이티브 고밀도만.
            # cover/section/visual/structural은 이 분기를 통과하지 않아 기존 경로(풀블리드/
            # NATIVE_SHAPES backdrop 등)를 그대로 유지한다. flag OFF에서는 완전한 no-op이라
            # 산출물이 바이트 단위로 동일하다(R6.1/6.4). content 편집 렌더는 풀블리드가 없으므로
            # _strip_text_over_fullbleed 대상이 아니며 편집 텍스트를 스트립하지 않는다(R4.6).
            if (_hybrid_on and not _native_routed
                    and not _rr_caller_img and not _rr_caller_bg
                    and not sd.get("nativeDiagram")):
                try:
                    _hc_role = _classify_slide_role(sd, False, title)
                except Exception:
                    _hc_role = "content"
                if _hc_role == "content":
                    # 데이터/토큰/팔레트를 기존 네이티브 경로와 동일 방식으로 구성한다.
                    try:
                        from ai_engine.slide_templates import design_tokens_for_profile as _dtfp_hc
                    except ImportError:
                        from slide_templates import design_tokens_for_profile as _dtfp_hc
                    try:
                        _hc_tokens = _dtfp_hc(style_profile)
                    except Exception:
                        _hc_tokens = {}
                    try:
                        _hc_palette = _build_palette(style_profile)
                    except Exception:
                        _hc_palette = None
                    _hc_bullets = [str(_b).strip() for _b in (bullets or []) if str(_b).strip()]
                    _hc_data = {"title": str(sd.get("title", "") or ""), "bullets": _hc_bullets}
                    # 히어로 = 이 슬라이드의 Vertex 사전생성 이미지(_vertex_pre[i]) 있으면 사용.
                    _hc_hero = _vertex_pre.get(i, "") or ""
                    try:
                        _hc_res = _render_content_editable(
                            s, prs, _hc_data, _hc_tokens, _hc_hero, _hc_palette) or {}
                        # 편집 경로로 라우팅 확정 — HTML 바이크/레거시 본문/후속 Vertex
                        # 풀블리드 임베드를 모두 스킵한다(_render_content_editable가 편집
                        # 가능 run>=1을 보장하므로 콘텐츠 손실 0).
                        _native_routed = True
                        _hybrid_content_routed = True
                        _hybrid_content_slide_ids.add(id(s))
                        _rr_hybrid_primary = "NATIVE_EDITABLE"
                        _rr_hybrid_slot = "hero" if _hc_hero else "none"
                        _rr_hybrid_editable = True
                        # task 8.5 — native 렌더러가 자체 제목 TextBox 를 방출했으면
                        # (title_count>=1) 루프 상단 _safe_set_title 이 채운 placeholder
                        # 제목을 제거해 제목 이중 방출(겹침)을 없앤다. 제목 방출 실패 시엔
                        # placeholder 제목을 보존한다(편집 제목 run>=1 유지, 손실-0).
                        try:
                            if int(_hc_res.get("title_count", 0) or 0) >= 1:
                                _hc_tph = s.shapes.title
                                if _hc_tph is not None:
                                    _hc_tph._element.getparent().remove(_hc_tph._element)
                        except Exception:
                            pass
                        print(f"[generate_pptx] 슬라이드 {i + 2} -> 하이브리드 content 편집 경로"
                              f"(layout={_hc_res.get('layout')}, "
                              f"image_placed={_hc_res.get('image_placed')}, "
                              f"image_preserved={_hc_res.get('image_preserved')})")
                    except Exception as _hce:
                        # content 편집 렌더 예외 — 콘텐츠 손실 0을 위해 기존 경로로 폴백(no-op).
                        print(f"[generate_pptx] 하이브리드 content 편집 실패(기존 경로 폴백) "
                              f"slide {i + 2}: {str(_hce)[:160]}")
            if (_html_enabled and not _native_routed
                    and not sd.get("slideBackground") and not sd.get("imageFile")
                    and not sd.get("nativeDiagram")):
                # 본문 주 렌더러 = HTML 고밀도 레이아웃(task 3.4). 성공 시 slideBackground가
                # 설정돼 네이티브 다이어그램 추론이 자동 skip된다. HTML이 렌더에 실패하면
                # (_sec_rel 빈 값) 아래 카드 폴백으로 본문 밀도를 유지한다(우선순위: HTML > 카드).
                try:
                    _sec_body = "\n".join(str(_b) for _b in bullets)
                    # task 3.5 — 이 슬라이드의 Vertex 사전생성 이미지를 HTML 렌더에
                    # 전달해 레이아웃 이미지 슬롯에 합성한다. 합성되면 단일 PNG에
                    # 이미지가 포함되고, 합성 불가 레이아웃이면 아래에서 on-slide
                    # 레이어링으로 폴백한다(둘 다 이미지를 폐기하지 않음 — 손실 0).
                    _pre_for_html = _vertex_pre.get(i, "")
                    _html_render_info = {}
                    _sec_rel = await _generate_html_slide_for_section(
                        _html_gw, _html_model, str(sd.get("title", "")), _sec_body, title, project_path,
                        style_profile=style_profile, hero_image=_pre_for_html, render_info=_html_render_info)
                    if _sec_rel:
                        if _pre_for_html and not _html_render_info.get("composited"):
                            # 합성 불가 HTML 레이아웃 + Vertex 이미지 → HTML PNG를 배경으로
                            # 쓰지 않고 on-slide 레이어링(이미지 back-most + 네이티브 콘텐츠)으로
                            # 폴백해 이미지를 보존한다. slideBackground 미설정 시 아래 카드
                            # 폴백 + _select_render_plan backdrop 경로가 pre를 풀블리드로 䧀다.
                            print(f"[generate_pptx] 슬라이드 {i + 2} HTML 레이아웃 이미지 슬롯 부재 → on-slide 레이어링 폴백")
                        else:
                            sd["slideBackground"] = _sec_rel
                            _rr_html_bg_set = True
                            if _pre_for_html and _html_render_info.get("composited"):
                                _pre_composited_into_html = True
                                print(f"[generate_pptx] 슬라이드 {i + 2} HTML+Vertex 합성(layout={_html_render_info.get('layout')}) — 단일 PNG")
                except Exception as _she:
                    print(f"[generate_pptx] 슬라이드 HTML 실패(카드 폴백): {str(_she)[:160]}")
                # 카드 폴백은 HTML 렌더 실패 시에만 동작(HTML 비활성 경로는 위에서 처리됨).
                if (not sd.get("slideBackground") and not sd.get("imageFile")
                        and not sd.get("nativeDiagram")):
                    try:
                        _fb = _cards_fallback_from_bullets(bullets)
                        if _fb:
                            sd["nativeDiagram"] = _fb
                            print(f"[generate_pptx] 슬라이드 {i + 2} HTML 미적용 → 카드 폴백(type={_fb.get('type')})")
                    except Exception as _fbe:
                        print(f"[generate_pptx] 카드 폴백 실패 slide {i + 2}: {str(_fbe)[:120]}")
            # 요구사항 9.5 — body placeholder 접근/채우기는 템플릿 단계로 간주해 격리한다.
            # 실패해도 슬라이드 자체는 유지되며 이후 슬라이드 생성이 계속된다(요구사항 9.6).
            # task 3.2 — HTML 베이크 후보였으나 베이크가 실패(_rr_html_bg_set False)해
            # 콘텐츠가 구워진 배경이 없으면, 제목을 네이티브로 폴백 설정한다(제목 손실 방지).
            if _html_bake_eligible and not _rr_html_bg_set and not _native_routed:
                _safe_set_title(s, sd.get("title", f"Slide {i + 2}"))
            elif _html_bake_eligible and (_rr_html_bg_set or _native_routed):
                # 콘텐츠가 구워진 배경 채택 — 도너에서 물려받은 제목 placeholder 텍스트를
                # 비워 _remove_empty_placeholders 가 제거하도록 한다(제목 1회: 베이크 PNG만).
                try:
                    _tph = s.shapes.title
                    if _tph is not None:
                        _tph.text = ""
                except Exception:
                    pass
            # task 3.1 — 콘텐츠가 구워진 풀블리드 배경(_rr_html_bg_set)이 채택된 슬라이드는
            # 그 PNG가 본문을 이미 포함하므로 네이티브 본문 텍스트박스를 방출하지 않는다
            # (중복 100% 겹침 제거). 그 외(비결함) 경로는 기존 동작을 그대로 보존한다.
            # task 9.1 — 네이티브 라우팅 슬라이드는 render_native_layout 가 본문을
            # 이미 편집가능 도형으로 방출했으므로 레거시 불릿 텍스트박스를 억제한다(중복 0).
            # 베이크 경로(_rr_html_bg_set)는 기존 통짜이미지 억제 동작을 그대로 보존(가산적).
            _suppress_native_body = bool(_rr_html_bg_set) or _native_routed
            body_shape = None
            try:
                body_shape = s.placeholders[1] if len(s.placeholders) > 1 else None
                if _suppress_native_body and body_shape is not None:
                    # 베이크 배경 채택 — 도너 상속 본문 텍스트를 비워 빈 placeholder 제거로
                    # 정리되게 한다(본문 중복 0). 비결함 경로는 _suppress False라 영향 없음.
                    try:
                        body_shape.text_frame.clear()
                    except Exception:
                        pass
                    body_shape = None
                if body_shape is None and bullets and not _suppress_native_body:
                    # 도너를 배경만 남기고 정리해 본문 placeholder가 없을 수 있다.
                    # 깔끔한 본문 텍스트박스를 생성해 불릿을 담는다(정렬 보장).
                    body_shape = s.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(12.1), Inches(5.4))
                if body_shape and bullets and not _suppress_native_body:
                    tf = body_shape.text_frame
                    tf.clear()
                    for j, bullet in enumerate(bullets):
                        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                        _add_md_paragraph(p, bullet)
            except Exception as e:
                print(f"[generate_pptx] body placeholder 채우기 실패 slide {i + 2}: {str(e)[:200]}")
                body_shape = None

            # Auto-generate image if imagePrompt is set, or use pre-rendered imageFile
            img_prompt = sd.get("imagePrompt", "")
            img_file = sd.get("imageFile", "")
            slide_bg = sd.get("slideBackground", "")
            native_diag = sd.get("nativeDiagram") if isinstance(sd, dict) else None
            # === 직접 호출 경로(모델→generate_pptx)도 편집 가능 다이어그램으로 ===
            # 모델이 nativeDiagram 없이 imagePrompt/bullets만 준 경우, 내용이 다이어그램형
            # (흐름/구조/트리/아키텍처)이면 통짜 PNG 대신 네이티브 도형으로 그린다.
            # 분류기가 빈 kind를 내는 순수 사진/일러스트 프롬프트는 기존 PNG 경로 유지.
            # === Vertex(Nano Banana Pro) 이미지 — 루프 진입 전 병렬 생성 결과 사용 ===
            # (순차 await로 인한 장시간 블로킹/끊김 방지). 결과가 있으면 네이티브 박스
            # 다이어그램 대신 고품질 이미지를 임베드한다.
            # 1) 네이티브(편집 가능) 다이어그램을 *먼저* 추론 — 구조형 슬라이드는
            # 정렬·편집 가능한 네이티브 도형으로. (이전: Vertex 래스터가 먼저 채워
            # 구조 슬라이드까지 통짜 AI 이미지가 되며 깨짐/오와열 발생.)
            if (not native_diag and not img_file and not slide_bg
                    and not _native_routed  # task20(D): 네이티브 라우팅 시 이중 렌더 차단
                    and os.environ.get("AE_DISABLE_NATIVE_DIAGRAM", "") != "1"):
                try:
                    _heading = str(sd.get("title", "") or "")
                    _body = "\n".join(str(b) for b in bullets)
                    _kind, _content = _classify_section_diagram(_heading, _body, title)
                    if not _kind and img_prompt:
                        _kind, _content = _classify_section_diagram(_heading, str(img_prompt), title)
                    if _kind:
                        native_diag = {"type": _kind, "content": (_content or _body or str(img_prompt))}
                        print(f"[generate_pptx] 슬라이드 {i + 2} → 네이티브 다이어그램(type={_kind})")
                except Exception as _ce:
                    print(f"[generate_pptx] 다이어그램 분류 실패 slide {i + 2}(무시): {str(_ce)[:160]}")
            # 2) 손실 없는 결정 규칙 — selectRenderPlan(task 3.3). 폐기형 가드
            #    (`if not native_diag and not img_file and not slide_bg: img_file = pre`)를
            #    역할 기반 손실-0 분기로 교체한다(design Fix Implementation §3).
            #    핵심 불변식: 생성된 Vertex 이미지(_pre_rel)는 어떤 분기에서도 폐기되지 않는다.
            # task 3.5(defect B): body_separated 신호를 풀블리드 배경 합성 경로로 전달.
            # 기본 False=기존 동작 보존(바이트 보존). 신호가 참일 때만 body_safe_area가
            # 본문을 배경과 분리된 안전 영역으로 이동/축소한다(additive 연결).
            _body_separated_signal = False
            _pre_rel = _vertex_pre.get(i)
            # task 7.4 — 하이브리드 content 편집 경로로 이미 히어로를 바운디드로 합성/보존한
            # 슬라이드는 여기서 풀블리드(visual/backdrop)로 재임베드하지 않는다(풀블리드 0 보존).
            if _pre_rel and not _pre_composited_into_html and not _hybrid_content_routed:
                try:
                    _slide_role = _classify_slide_role(sd, False, title)
                except Exception:
                    _slide_role = "content"
                _plan = _select_render_plan(
                    has_vertex_image=True,
                    has_native_diagram=bool(native_diag),
                    has_image_file=bool(img_file),
                    has_slide_bg=bool(slide_bg),
                    role=_slide_role,
                    html_enabled=bool(_html_enabled),
                )
                _slot = _plan["vertex_slot"]
                _body_separated_signal = bool(_plan.get("body_separated"))
                _rr_role = _slide_role
                _rr_plan_slot = _slot
                # === task 7.2: 역할 기반 주 렌더러 선택(플랜 결정만, 가산적·게이팅) ===
                # `_hybrid_on == True`일 때만, 그리고 caller가 imageFile/slideBackground를
                # 지정하지 않은 슬라이드(Vertex 사전생성이 caller 지정 슬라이드를 이미
                # skip하므로 이 seam에 도달하는 슬라이드는 본질적으로 caller 미지정)에 한해,
                # 위에서 예외 안전하게 확정된 역할(_slide_role, 예외 시 "content")로
                # _select_hybrid_render_plan을 호출해 주 렌더러를 결정한다. 본 태스크는 플랜
                # *결정*과 관측 기록만 수행하며 실제 렌더 경로 스왑(프롬프트 7.3 / content
                # 편집 7.4)은 별도 태스크다. 따라서 렌더 출력은 변하지 않고, flag OFF에서는
                # 이 분기가 전혀 실행되지 않아 산출물이 바이트 단위로 동일하다(R6.1/6.4).
                # caller 지정 미디어 우선순위(R4.2)와 손실-0 최종 검증
                # (has_vertex_image ⇒ slot≠none, R4.1)은 위 _select_render_plan에 위임 유지.
                if _hybrid_on and not _rr_caller_img and not _rr_caller_bg:
                    _hybrid_plan = _select_hybrid_render_plan(
                        role=_slide_role,
                        vertex_enabled=bool(_rr_vertex_enabled),
                        html_enabled=bool(_html_enabled),
                        has_vertex_image=True,
                        has_native_diagram=bool(native_diag),
                        has_image_file=bool(img_file),
                        has_slide_bg=bool(slide_bg),
                    )
                    _rr_hybrid_primary = _hybrid_plan["primary"]
                    _rr_hybrid_slot = _hybrid_plan["vertex_slot"]
                    _rr_hybrid_editable = bool(_hybrid_plan["editable"])
                    print(f"[generate_pptx] 슬라이드 {i + 2} -> 하이브리드 플랜 결정"
                          f"(role={_slide_role}, primary={_rr_hybrid_primary}, "
                          f"slot={_rr_hybrid_slot}, editable={_rr_hybrid_editable})")
                if _slot == "visual":
                    # 비구조 비주얼(cover/visual/content) — 이미지가 곧 슬라이드 비주얼.
                    # (기존 가드 동작 보존: not native_diag and not img_file and not slide_bg)
                    img_file = _pre_rel
                    print(f"[generate_pptx] 슬라이드 {i + 2} -> Vertex 이미지 임베드(role={_slide_role}, primary=image)")
                elif _slot == "backdrop":
                    # 네이티브 도형(구조형/카드) 또는 caller 이미지 유지 + pre를 풀블리드
                    # backdrop으로 보존(폐기 금지). slide_bg로 두면 아래 _eff_bg 경로가
                    # back-most로 깔고 그 위에 네이티브 도형을 올린다(_native_over_bg).
                    if not slide_bg:
                        slide_bg = _pre_rel
                        print(f"[generate_pptx] 슬라이드 {i + 2} -> Vertex 이미지 backdrop 보존(role={_slide_role}, primary={_plan['primary']})")
                elif _slot == "hero":
                    # HTML 풀블리드가 주 렌더러. pre는 히어로/이미지 슬롯 합성용으로 보존한다
                    # (task 3.5에서 slide_templates 이미지 슬롯에 주입). 합성 전에도 폐기하지
                    # 않도록 heroImage 힌트로 기록(손실 0).
                    sd["heroImage"] = _pre_rel
                    print(f"[generate_pptx] 슬라이드 {i + 2} -> Vertex 이미지 히어로 슬롯 보존(HTML 주 렌더러)")
            img_path = ""
            bg_path = ""
            native_drawn = False
            # 문제2 — 네이티브(편집 가능) 다이어그램 우선 처리.
            # nativeDiagram spec({type, content})이 있으면 통짜 PNG 대신 python-pptx
            # 도형으로 우측 절반에 직접 조립한다. 본문 placeholder는 좌측 절반으로 리사이즈.
            # 실패 시 native_drawn=False → 아래 imageFile(PNG) 폴백 경로로 진행.
            # 본문도 고품질: slideBackground(Vertex 배경)가 있으면 먼저 풀블리드로 깔고,
            # 그 위에 흰 콘텐츠 카드+편집 가능 다이어그램을 그린다(품질감+편집성 동시).
            _native_over_bg = False
            # D1(슬라이드당 풀블리드 1회 보장): 이 슬라이드에 풀블리드 배경이 이미
            # 임베드됐는지 추적하는 슬라이드 단위 가드 플래그(슬라이드마다 False 초기화).
            # 결정은 순수 함수 layout_geometry.fullbleed_guard 에 위임한다(PBT 대상).
            _fb_embedded = False
            try:
                from layout_geometry import fullbleed_guard as _fb_guard
            except Exception:
                try:
                    from ai_engine.layout_geometry import fullbleed_guard as _fb_guard
                except Exception:
                    _fb_guard = lambda _c: _c < 1
            _eff_bg = slide_bg or _dp_body_bg
            if native_diag and isinstance(native_diag, dict) and _eff_bg and _fb_guard(1 if _fb_embedded else 0):
                try:
                    _cand_bg = _eff_bg if os.path.isabs(_eff_bg) else os.path.join(_local_root, _eff_bg)
                    if os.path.isfile(_cand_bg):
                        _pic_bg = s.shapes.add_picture(_cand_bg, Inches(0), Inches(0),
                                                       width=Inches(13.333), height=Inches(7.5))
                        try:
                            _spt = _pic_bg._element.getparent()
                            _spt.remove(_pic_bg._element)
                            _spt.insert(2, _pic_bg._element)
                        except Exception:
                            pass
                        _native_over_bg = True
                        _fb_embedded = True  # D1: 이 슬라이드 풀블리드 1회 소모
                        slide_bg = ""
                except Exception as _ebg:
                    print(f"[generate_pptx] \ubcf8\ubb38 \ubc30\uacbd \uc784\ubca0\ub4dc \uc2e4\ud328 slide {i + 2}: {str(_ebg)[:160]}")
            # B방향(수읅2): 통짜 풀블리드 배경(_native_over_bg)이 채택된 다이어그램
            # 슬라이드는 그 배경만 남긴다 — build_native_diagram 네이티브 카드/텍스트/
            # 도형 방출을 억제하고(공존 금지) 본문/제목 등 겹치는 편집 셀입를 제거한다
            # (겹츨0·중복0). 배경(HTML 베이크/Vertex)이 콘텐츠를 포함한다고 신뢰. 손실0:
            # 임베드된 배경 이미지는 폐기하지 않는다.
            if _native_over_bg:
                native_drawn = True  # 이미 통짜 배경 임베드 → 아래 이미지/배경 폴백 skip
                try:
                    _strip_text_over_fullbleed(s)
                except Exception:
                    pass
            if native_diag and isinstance(native_diag, dict) and not _native_over_bg:
                try:
                    from native_diagram_pptx import build_native_diagram as _bnd
                except Exception:
                    try:
                        from ai_engine.native_diagram_pptx import build_native_diagram as _bnd
                    except Exception:
                        _bnd = None
                if _bnd is not None:
                    try:
                        from pptx.util import Emu as _Emu2
                        # 다이어그램은 제목 아래 전체 폭을 사용한다(이전: 우측 절반 6인치에
                        # 눌려 박스·텍스트가 cramped/정렬 어긋남). 동일 내용의 bullets는
                        # 다이어그램이 대체하므로 본문 placeholder를 비워 중복·간섭 제거.
                        # task 3.5(defect B): 본문 region 을 배경 이미지와 분리.
                        # _native_over_bg 로 풀블리드 배경을 깐 경우 bg_rect 를 넘긴다.
                        # bg 가 백드롭(흰 패널 분리, has_baked_text=False)이면 desired
                        # 보존(바이트 보존). 구워진-텍스트 신호가 참일 때만 안전 영역으로
                        # 이동/축소. 손실-0: 생성된 Vertex 이미지는 폐기되지 않는다.
                        _body_region = (0.6, 1.7, 12.1, 5.2)
                        try:
                            try:
                                from layout_geometry import body_safe_area as _bsa3
                            except Exception:
                                from ai_engine.layout_geometry import body_safe_area as _bsa3
                            _bg_rect3 = (0.0, 0.0, 13.333, 7.5) if _native_over_bg else None
                            _body_region = _bsa3(
                                (0.0, 0.0, 13.333, 7.5), _bg_rect3,
                                has_baked_text=_body_separated_signal,
                                desired=_body_region,
                            )
                        except Exception:
                            _body_region = (0.6, 1.7, 12.1, 5.2)
                        native_drawn = _bnd(
                            s,
                            (native_diag.get("type") or "block"),
                            (native_diag.get("content") or ""),
                            region=_body_region,
                            palette=_tpl_palette_for_native,
                            title="",
                            backdrop=_native_over_bg,  # 제목은 슬라이드 placeholder가 표시 — 다이어그램 내 중복 방지(이슈3)
                            note=(native_diag.get("note") or ""),
                        )
                        if native_drawn and body_shape is not None:
                            # 본문 텍스트 비우기(다이어그램이 내용을 표현) — 빈 placeholder는
                            # 화면에 보이지 않으므로 슬라이드가 깔끔해진다.
                            try:
                                body_shape.text_frame.clear()
                            except Exception:
                                pass
                        if native_drawn:
                            try:
                                _rb = _remove_content_band_pictures(s, 1.7, 6.9)
                                if _rb:
                                    print(f"[generate_pptx] 다이어그램 영역 겹침 장식 띠 {_rb}개 제거 slide {i + 2}")
                            except Exception:
                                pass
                            print(f"[generate_pptx] 네이티브 다이어그램(전체폭) slide {i + 2} (type={native_diag.get('type')})")
                    except Exception as _nde:
                        print(f"[generate_pptx] 네이티브 다이어그램 실패 slide {i + 2} → PNG 폴백: {str(_nde)[:200]}")
                        native_drawn = False
            try:
                if native_drawn:
                    # 네이티브로 그렸으면 PNG/배경/이미지 단계를 모두 건너뛴다.
                    # (이전 버그: 아래 `if img_file`/`elif img_prompt`가 별도 if라서
                    #  native_drawn=True여도 래스터 이미지가 추가돼 도형과 겹쳤다.)
                    pass
                else:
                    if slide_bg:
                        # Full-bleed background path (HTML→PNG capture).
                        cand = slide_bg if os.path.isabs(slide_bg) else os.path.join(_local_root, slide_bg)
                        if os.path.isfile(cand):
                            bg_path = cand
                        else:
                            print(f"[generate_pptx] slideBackground not found: {slide_bg}")
                    if img_file and not bg_path:
                        cand = img_file if os.path.isabs(img_file) else os.path.join(_local_root, img_file)
                        if os.path.isfile(cand):
                            img_path = cand
                        else:
                            print(f"[generate_pptx] imageFile not found: {img_file}")
                    elif img_prompt and not bg_path:
                        img_result_str = await _tool_generate_image({"prompt": img_prompt, "size": "1024x1024"},
                            project_path, aws_profile=aws_profile, bedrock_user=bedrock_user)  # [patched-credentials]
                        img_result = json.loads(img_result_str)
                        if "path" in img_result:
                            rel = img_result["path"]
                            cand = rel if os.path.isabs(rel) else os.path.join(_local_root, rel)
                            if os.path.isfile(cand):
                                img_path = cand
            except Exception as e:
                # 이미지 단계 실패가 outer try/except로 새지 않도록 차단 — 슬라이드는 텍스트만 유지
                print(f"[generate_pptx] image step failed slide {i + 2}: {e}")
                img_path = ""
                bg_path = ""

            if bg_path and _fb_guard(1 if _fb_embedded else 0):
                try:
                    # 요구사항 6.5 — Full-bleed slide background (1920×1080 HTML capture)
                    # 가 전체 13.333×7.5in 슬라이드를 덮는다. 그림을 (0,0)에 추가한 뒤
                    # spTree.insert(2, ...)로 첫 shape 자리(back-most)에 두면:
                    #   - 슬라이드 마스터 배경보다 위 레이어 → 템플릿 마스터 배경을 가린다.
                    #   - add_slide가 먼저 복제한 title/body placeholder들은 인덱스 3+에 남아
                    #     그림 위에 렌더 → 텍스트는 항상 위, 편집 가능 유지(요구사항 6.8).
                    # python-pptx는 z-order를 직접 노출하지 않으므로 spTree 인덱스 조작으로
                    # 처리한다(0=nvGrpSpPr, 1=grpSpPr, 2+=shapes). 템플릿/무템플릿 동일.
                    pic = s.shapes.add_picture(
                        bg_path,
                        Inches(0), Inches(0),
                        width=Inches(13.333), height=Inches(7.5),
                    )
                    # 마스터 배경 위·placeholder 아래로 이동(요구사항 6.5).
                    try:
                        spTree = pic._element.getparent()
                        spTree.remove(pic._element)
                        spTree.insert(2, pic._element)  # 0=nvGrpSpPr, 1=grpSpPr, 2+=shapes
                    except Exception:
                        pass
                    _fb_embedded = True  # D1: 이 슬라이드 풀블리드 1회 소모
                except Exception as e:
                    print(f"[generate_pptx] full-bleed embed failed slide {i + 2}: {e}")
            elif img_path:
                try:
                    # 이미지 비율을 측정해 영역에 fit(비율 보존) + 중앙정렬.
                    # 이전 버그: width=6.0만 지정 → 정사각/세로 이미지가 슬라이드 하단(7.5in)을
                    # 넘쳐 잘리고, 텍스트가 적은 슬라이드에서도 우측 절반에만 작게 박혔다.
                    from pptx.util import Emu as _Emu
                    iw, ih = 4, 3
                    try:
                        from PIL import Image as _PILfit
                        with _PILfit.open(img_path) as _im:
                            iw, ih = _im.size
                    except Exception:
                        pass
                    ar = (iw / ih) if ih else 1.3333

                    _has_text = bool(bullets) and any(str(b).strip() for b in bullets)
                    if _has_text:
                        # 2-칼럼: 본문 좌측, 이미지는 우측 영역에 fit.
                        if body_shape is not None:
                            try:
                                body_shape.left = _Emu(int(0.6 * 914400))
                                body_shape.top = _Emu(int(1.6 * 914400))
                                body_shape.width = _Emu(int(6.0 * 914400))
                                body_shape.height = _Emu(int(5.4 * 914400))
                            except Exception as _le:
                                print(f"[generate_pptx] body resize 실패 slide {i + 2}: {_le}")
                        region_l, region_t, region_w, region_h = 7.0, 1.6, 6.0, 5.4
                    else:
                        # 이미지 중심 슬라이드 — 크게 중앙 배치(작게 박히는 문제 해결).
                        region_l, region_t, region_w, region_h = 1.5, 1.7, 10.33, 5.2

                    # D2(슬롯-이미지 크기 정합): 대상 region 이 소형 장식 슬롯인데 대형
                    # 이미지가 흘러들면 콘텐츠 region 으로 승격한다(손실-0: 이미지 폐기 금지).
                    try:
                        try:
                            from layout_geometry import slot_image_fits as _sif
                        except Exception:
                            from ai_engine.layout_geometry import slot_image_fits as _sif
                        if not _sif((region_l, region_t, region_w, region_h), iw, ih):
                            region_l, region_t, region_w, region_h = 1.5, 1.7, 10.33, 5.2
                    except Exception:
                        pass
                    # 영역에 비율 보존 fit
                    draw_w = region_w
                    draw_h = region_w / ar if ar else region_h
                    if draw_h > region_h:
                        draw_h = region_h
                        draw_w = region_h * ar
                    off_l = region_l + (region_w - draw_w) / 2.0
                    off_t = region_t + (region_h - draw_h) / 2.0
                    # D3(경계 클램프): 최종 배치 rect 를 슬라이드 경계 안으로 클램프해
                    # 음수 top/left·경계 초과를 모든 region 정의에서 제거한다(이미 경계 안이면
                    # no-op → 바이트 보존). 손실-0: 이미지는 폐기되지 않고 경계 안으로 이동.
                    try:
                        try:
                            from layout_geometry import clamp_into_bounds as _clampb
                        except Exception:
                            from ai_engine.layout_geometry import clamp_into_bounds as _clampb
                        off_l, off_t, draw_w, draw_h = _clampb((off_l, off_t, draw_w, draw_h))
                    except Exception:
                        pass
                    s.shapes.add_picture(
                        img_path,
                        Inches(off_l), Inches(off_t),
                        width=Inches(draw_w), height=Inches(draw_h),
                    )
                except Exception as e:
                    print(f"[generate_pptx] embed failed slide {i + 2}: {e}")

            # 슬라이드 합성 종료 — 채우지 않은 빈 placeholder(프롬프트 "텍스트를 입력하십시오"/
            # 점선 테두리) 제거. 네이티브 다이어그램으로 본문을 비운 경우, bullets/이미지가
            # 없어 본문이 빈 경우 모두 깔끔하게 정리된다.
            _remove_empty_placeholders(s)
            _remove_empty_text_shapes(s)
            _clamp_shapes_into_bounds(s)  # B방향(수읅5): 도형 경계 후처리 클램프
            # === renderReport — 이 슬라이드의 실제 렌더 경로 기록(관측만) ===
            try:
                _rr_role_final = _rr_role
                if not _rr_role_final:
                    try:
                        _rr_role_final = _classify_slide_role(sd, False, title)
                    except Exception:
                        _rr_role_final = "content"
                _rr_pre_for_slide = _vertex_pre.get(i)
                _rr_vembed = bool(_pre_composited_into_html
                                  or (_hybrid_content_routed and _rr_hybrid_slot == "hero")
                                  or (_rr_pre_for_slide and _rr_plan_slot in ("visual", "backdrop", "hero")))
                if _hybrid_content_routed:
                    _rr_spath = "content-editable"
                elif _pre_composited_into_html or _rr_html_bg_set:
                    _rr_spath = "html-fullbleed"
                elif _rr_plan_slot == "visual":
                    _rr_spath = "vertex-visual"
                elif _rr_plan_slot == "hero":
                    _rr_spath = "vertex-hero"
                elif _rr_plan_slot == "backdrop":
                    if native_drawn:
                        _rr_spath = "native-backdrop"
                    elif img_path and _rr_caller_img:
                        _rr_spath = "caller-image"
                    else:
                        _rr_spath = "native-backdrop"
                elif native_drawn:
                    _rr_spath = "native-shapes"
                elif _rr_caller_img and img_path:
                    _rr_spath = "caller-image"
                elif bg_path:
                    _rr_spath = "html-fullbleed"
                elif img_path:
                    _rr_spath = "caller-image"
                else:
                    _rr_spath = "text-only"
                _rr_slides.append({"index": i + 1, "role": _rr_role_final,
                                   "path": _rr_spath, "vertexEmbedded": _rr_vembed})
            except Exception as _rrse:
                print(f"[generate_pptx] renderReport slide 기록 실패(무시): {str(_rrse)[:120]}")

        # 템플릿 도너 중 새 내용으로 사용되지 않은 잉여 슬라이드 제거.
        # (콘텐츠 슬라이드 수 < 템플릿 디자인 슬라이드 수인 경우 — 남은 샘플 슬라이드 삭제)
        if _use_designed and _donor_consumed["n"] < len(_donors):
            for _unused in _donors[_donor_consumed["n"]:]:
                _remove_slide_obj(prs, _unused)

        # === 전 슬라이드 공통 디자인 베이스라인 (내용 무관 자동 미화) ===
        # 무템플릿(기본 레이아웃)이고 HTML 풀블리드 배경을 쓰지 않을 때만 적용.
        # 템플릿/HTML 경로는 자체 디자인이 있으므로 건드리지 않는다.
        try:
            if not _html_enabled:
                _apply_universal_slide_design(
                    prs,
                    palette=_tpl_palette_for_native,
                    doc_title=title,
                    used_template=used_template,
                    skip_slide_ids=_hybrid_content_slide_ids,
                )
        except Exception as _ud_e:
            print(f"[generate_pptx] universal design 적용 실패(무시): {str(_ud_e)[:160]}")

        # task 3.3 — 경계 밖 도형 클램프 후처리 패스 (defect 1.4).
        # 모든 슬라이드의 모든 도형 rect 를 layout_geometry.within_bounds 로 검사하고,
        # 슬라이드 경계(0,0,13.333,7.5) 밖이면 clamp_into_bounds 로 경계 안으로 보정한다.
        # 경계 안 도형은 within_bounds==True → clamp 가 입력을 그대로 반환(no-op)하므로
        # 비결함 슬라이드의 도형 좌표는 바이트 보존된다(additive).
        try:
            try:
                from layout_geometry import within_bounds as _wb_chk, clamp_into_bounds as _clamp_b
            except Exception:
                from ai_engine.layout_geometry import within_bounds as _wb_chk, clamp_into_bounds as _clamp_b
            from pptx.util import Emu as _EmuClamp
            _EMU_IN = 914400.0
            for _cl_slide in prs.slides:
                for _cl_sh in list(_cl_slide.shapes):
                    try:
                        _cl_l, _cl_t, _cl_w, _cl_h = _cl_sh.left, _cl_sh.top, _cl_sh.width, _cl_sh.height
                        if None in (_cl_l, _cl_t, _cl_w, _cl_h):
                            continue
                        _cl_rect = (_cl_l / _EMU_IN, _cl_t / _EMU_IN,
                                    _cl_w / _EMU_IN, _cl_h / _EMU_IN)
                        if _wb_chk(_cl_rect):
                            continue
                        _nl, _nt, _nw, _nh = _clamp_b(_cl_rect)
                        _cl_sh.left = _EmuClamp(int(round(_nl * _EMU_IN)))
                        _cl_sh.top = _EmuClamp(int(round(_nt * _EMU_IN)))
                        _cl_sh.width = _EmuClamp(int(round(_nw * _EMU_IN)))
                        _cl_sh.height = _EmuClamp(int(round(_nh * _EMU_IN)))
                    except Exception:
                        continue
        except Exception as _clamp_e:
            print(f"[generate_pptx] 경계 클램프 후처리 실패(무시): {str(_clamp_e)[:160]}")

        try:
            prs.save(output_path)
        except Exception as save_err:
            return json.dumps({"error": "pptx-generation-failed", "detail": str(save_err)[:200]})

        # TASK 8 — 응답 직전 디스크 자체 재검증.
        # prs.save가 예외 없이 끝나도 file-system race / cleanup hook / 권한 충돌로
        # 0바이트 또는 미존재 상태가 될 수 있음. 그 상태로 path를 응답하면 패널이
        # 표시한 뒤 사용자가 다운로드 시 sourcePath가 빈 파일/없음.
        if not os.path.isfile(output_path):
            return json.dumps({
                "error": "pptx-generation-failed",
                "detail": f"save reported success but file missing: {output_path}",
            })
        size_bytes = os.path.getsize(output_path)
        if size_bytes <= 0:
            try:
                os.remove(output_path)
            except OSError:
                pass
            return json.dumps({
                "error": "pptx-generation-failed",
                "detail": "save produced zero-byte file",
            })
        result = {
            "path": relative_path,
            # TASK 8 근본수정 — 도구가 실제 저장한 절대경로를 응답에 포함.
            # 프론트 채팅 카드(resolveFullPath)가 state.folderPath+상대경로로
            # 추측 조립하던 것을 absPath 우선 사용으로 교체 → 원격/패키징/폴더미오픈
            # 모든 케이스에서 다운로드·미리보기 경로가 실제 저장 위치와 일치.
            "absPath": output_path,
            "model": "python-pptx",
            "slideCount": len(slides_data) + 1,  # +1 for cover
            "sizeBytes": size_bytes,
        }
        # === renderReport (additive) — 선택된 렌더 경로 관측치를 결과에 노출.
        # 기존 키는 제거/변경하지 않고 새 키 하나만 추가한다(pptx 출력 바이트 불변).
        try:
            _rr_cover_entry = {"index": 0, "role": "cover",
                               "path": (_rr_cover_path or ("html-fullbleed" if _html_enabled else "native-shapes")),
                               "vertexEmbedded": False}
            _rr_all = [_rr_cover_entry] + list(_rr_slides)
            _rr_generated = len(_vertex_pre)
            _rr_embedded = sum(1 for _s in _rr_all if _s.get("vertexEmbedded"))
            _rr_unused = _rr_generated - _rr_embedded
            if _rr_unused < 0:
                _rr_unused = 0
            result["renderReport"] = {
                "htmlEnabled": bool(_html_enabled),
                "htmlRenderer": _rr_html_renderer,
                "htmlDisabledReason": _rr_html_reason,
                "vertexEnabled": bool(_rr_vertex_enabled),
                "vertexDisabledReason": _rr_vertex_reason,
                "slideCount": len(_rr_all),
                "slides": _rr_all,
                "vertexGenerated": _rr_generated,
                "vertexEmbedded": _rr_embedded,
                "vertexUnused": _rr_unused,
            }
        except Exception as _rre:
            print(f"[generate_pptx] renderReport 생성 실패(무시): {str(_rre)[:160]}")
        # 요구사항 6.7 — 템플릿 적용 성공 시에만 templateId를 응답에 포함.
        # 무템플릿 경로는 기존 응답 형태를 그대로 유지(하위 호환, 요구사항 5.2).
        if used_template and template_id:
            result["templateId"] = template_id
        return json.dumps(result)
    except Exception as e:
        return json.dumps({"error": "pptx-generation-failed", "detail": str(e)[:200]})


async def _tool_generate_xlsx(tool_input: dict, project_path: str) -> str:
    """Generate an Excel workbook (.xlsx) using openpyxl. Accepts lenient input shapes."""
    title, sheets_data = _normalize_doc_input(tool_input, default_kind="sheets")

    if not title:
        return json.dumps({"error": "title is required"})
    if not sheets_data:
        return json.dumps({"error": "sheets is required"})

    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "openpyxl", "hint": "pip install openpyxl"})

    import time as _t, re as _re
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _slug_from_title(title) or "workbook"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.xlsx"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    try:
        wb = Workbook()
        wb.remove(wb.active)  # remove default empty sheet
        header_font = Font(bold=True, color="FFFFFFFF")
        header_fill = PatternFill(start_color="FF007ACC", end_color="FF007ACC", fill_type="solid")
        header_align = Alignment(horizontal="center", vertical="center")
        thin = Side(border_style="thin", color="FFCCCCCC")
        border = Border(left=thin, right=thin, top=thin, bottom=thin)

        for i, sd in enumerate(sheets_data):
            if not isinstance(sd, dict):
                sd = {"name": f"Sheet{i+1}", "headers": [], "rows": []}
            sheet_name = (sd.get("name") or f"Sheet{i+1}")[:31]
            ws = wb.create_sheet(title=sheet_name)
            headers = sd.get("headers") or []
            rows = sd.get("rows") or []
            # Headers
            for c, h in enumerate(headers, start=1):
                cell = ws.cell(row=1, column=c, value=str(h))
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = header_align
                cell.border = border
            # Data rows
            for r, row in enumerate(rows, start=2):
                if not isinstance(row, list):
                    row = [row]
                for c, val in enumerate(row, start=1):
                    cell = ws.cell(row=r, column=c, value=val)
                    cell.border = border
            # Auto column widths (cap at 60 chars)
            for c in range(1, max(1, len(headers)) + 1):
                col_letter = get_column_letter(c)
                max_len = len(str(headers[c-1])) if c <= len(headers) else 8
                for row in rows:
                    if isinstance(row, list) and c <= len(row):
                        max_len = max(max_len, len(str(row[c-1])))
                ws.column_dimensions[col_letter].width = min(max_len + 2, 60)

        if not wb.sheetnames:
            # Empty input — create a placeholder sheet so the file is valid
            ws = wb.create_sheet(title="Sheet1")
            ws.cell(row=1, column=1, value=title)

        wb.save(output_path)
        # TASK 8 — save 후 디스크 자체 검증 (응답 path가 실재 보장).
        if not os.path.isfile(output_path):
            return json.dumps({
                "error": "xlsx-generation-failed",
                "detail": f"save reported success but file missing: {output_path}",
            })
        size_bytes = os.path.getsize(output_path)
        if size_bytes <= 0:
            try:
                os.remove(output_path)
            except OSError:
                pass
            return json.dumps({
                "error": "xlsx-generation-failed",
                "detail": "save produced zero-byte file",
            })
        return json.dumps({
            "path": relative_path,
            "absPath": output_path,  # TASK 8 근본수정 — 실제 저장 절대경로(카드 다운로드 정확성)
            "model": "openpyxl",
            "sheetCount": len(sheets_data),
            "sizeBytes": size_bytes,
        })
    except Exception as e:
        return json.dumps({"error": "xlsx-generation-failed", "detail": str(e)[:200]})


async def _tool_generate_docx(tool_input: dict, project_path: str, aws_profile: str = '', bedrock_user: str = '') -> str:  # [patched-credentials]
    """Generate a Word document (.docx) using python-docx. Accepts lenient input shapes.

    각 섹션은 선택적으로 imagePrompt를 가질 수 있으며, 존재 시
    _tool_generate_image로 이미지를 생성하고 heading 다음, body 이전에
    add_picture로 임베드한다 (실패 시 텍스트만으로 진행).
    """
    title, sections = _normalize_doc_input(tool_input, default_kind="sections")

    if not title:
        return json.dumps({"error": "title is required"})
    if not sections:
        return json.dumps({"error": "sections is required"})

    try:
        from docx import Document
        from docx.shared import Pt, Inches
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "python-docx", "hint": "pip install python-docx"})

    import time as _t, re as _re
    _local_root = _resolve_local_root(project_path)
    gen_dir = os.path.join(_local_root, ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _slug_from_title(title) or "doc"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.docx"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    try:
        doc = Document()
        # Default font tweak — readable size
        try:
            style = doc.styles["Normal"]
            style.font.size = Pt(11)
        except Exception:
            pass

        doc.add_heading(title, level=0)

        for sec in sections:
            if not isinstance(sec, dict):
                doc.add_paragraph(str(sec))
                continue
            heading = sec.get("heading", "")
            level = int(sec.get("level") or 2)
            level = max(1, min(level, 3))
            body = sec.get("body", "")
            bullets = sec.get("bullets") or []
            img_prompt = sec.get("imagePrompt", "")
            img_file = sec.get("imageFile", "")

            if heading:
                doc.add_heading(heading, level=level)

            # 섹션 이미지 — heading 다음, body 이전에 삽입 (fail-soft)
            # imageFile(사전 렌더) 우선, 없으면 imagePrompt(Bedrock) 사용
            img_abs = ""
            try:
                if img_file:
                    cand = img_file if os.path.isabs(img_file) else os.path.join(_local_root, img_file)
                    if os.path.isfile(cand):
                        img_abs = cand
                    else:
                        print(f"[generate_docx] imageFile not found: {img_file}")
                elif img_prompt:
                    img_result_str = await _tool_generate_image(
                        {"prompt": img_prompt, "size": "1024x1024"},
                        project_path, aws_profile=aws_profile, bedrock_user=bedrock_user,
                    )
                    img_result = json.loads(img_result_str)
                    if "path" in img_result:
                        rel = img_result["path"]
                        cand = rel if os.path.isabs(rel) else os.path.join(_local_root, rel)
                        if os.path.isfile(cand):
                            img_abs = cand
            except Exception as e:
                # 이미지 단계 실패가 outer try/except로 새지 않도록 차단 — 섹션은 텍스트 전용
                print(f"[generate_docx] image step failed section '{heading[:40]}': {e}")
                img_abs = ""

            if img_abs:
                try:
                    doc.add_picture(img_abs, width=Inches(6))
                except Exception as e:
                    print(f"[generate_docx] embed failed section '{heading[:40]}': {e}")

            if body:
                for para in str(body).split("\n"):
                    if para.strip():
                        doc.add_paragraph(para)
            for b in bullets:
                doc.add_paragraph(str(b), style="List Bullet")

        doc.save(output_path)
        # TASK 8 — save 후 디스크 자체 검증.
        if not os.path.isfile(output_path):
            return json.dumps({
                "error": "docx-generation-failed",
                "detail": f"save reported success but file missing: {output_path}",
            })
        size_bytes = os.path.getsize(output_path)
        if size_bytes <= 0:
            try:
                os.remove(output_path)
            except OSError:
                pass
            return json.dumps({
                "error": "docx-generation-failed",
                "detail": "save produced zero-byte file",
            })
        # python-docx exposes paragraphs via doc.paragraphs
        para_count = len(doc.paragraphs)
        return json.dumps({
            "path": relative_path,
            "absPath": output_path,  # TASK 8 근본수정 — 실제 저장 절대경로(카드 다운로드 정확성)
            "model": "python-docx",
            "sectionCount": len(sections),
            "paragraphCount": para_count,
            "sizeBytes": size_bytes,
        })
    except Exception as e:
        return json.dumps({"error": "docx-generation-failed", "detail": str(e)[:200]})


async def _tool_edit_image(tool_input: dict, project_path: str, aws_profile: str = '', bedrock_user: str = '') -> str:  # [patched-credentials]
    """Edit an image using inpaint or outpaint."""
    import time as _t, base64
    mode = tool_input.get("mode", "")
    image_path = tool_input.get("image_path", "")
    prompt = (tool_input.get("prompt") or "").strip()

    if mode not in ("inpaint", "outpaint"):
        return json.dumps({"error": "invalid-mode", "detail": "mode must be inpaint or outpaint"})
    if not image_path:
        return json.dumps({"error": "invalid-parameter", "detail": "image_path is required"})
    if not prompt:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt is required"})
    if len(prompt) > 512:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt exceeds 512 chars"})

    # Resolve path (relative to project_path)
    if not os.path.isabs(image_path) and project_path:
        full_path = os.path.join(project_path, image_path)
    else:
        full_path = image_path
    if not os.path.isfile(full_path):
        return json.dumps({"error": "file-not-found", "detail": f"image not found: {image_path}"})

    # Validate format and size
    # Req 3.5: outpaint surfaces "invalid-input"; inpaint (Req 2.9) keeps "invalid-image"
    _fmt_err = "invalid-input" if mode == "outpaint" else "invalid-image"
    try:
        with open(full_path, "rb") as f:
            magic = f.read(8)
        if magic[:8] == b"\x89PNG\r\n\x1a\n":
            fmt = "png"
        elif magic[:3] == b"\xff\xd8\xff":
            fmt = "jpeg"
        elif magic[:4] == b"RIFF":
            fmt = "webp"
        else:
            return json.dumps({"error": _fmt_err, "detail": "unsupported format (PNG/JPEG/WEBP only)"})
    except Exception as e:
        return json.dumps({"error": _fmt_err, "detail": str(e)[:200]})

    file_size = os.path.getsize(full_path)
    if file_size > 5 * 1024 * 1024:
        return json.dumps({"error": "invalid-image", "detail": "image exceeds 5MB"})

    # Req 2.9: inpaint allows only PNG/JPEG (outpaint additionally accepts WEBP per Req 3.5)
    if mode == "inpaint" and fmt == "webp":
        return json.dumps({"error": "invalid-image", "detail": "inpaint requires PNG or JPEG"})

    # Encode image
    with open(full_path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode("ascii")

    # Build request body per mode
    # [patched-credentials] honor explicit kw-args, fall back to env
    aws_profile = aws_profile or os.environ.get("AWS_PROFILE", "bedrock-gw")
    bedrock_user = bedrock_user or os.environ.get("BEDROCK_USER", "")
    gw = _get_gw(aws_profile, bedrock_user)

    last_error = ""

    if mode == "inpaint":
        mask_path = tool_input.get("mask_path", "")
        if not mask_path:
            return json.dumps({"error": "invalid-parameter", "detail": "mask_path required for inpaint"})
        if not os.path.isabs(mask_path) and project_path:
            mask_full = os.path.join(project_path, mask_path)
        else:
            mask_full = mask_path
        if not os.path.isfile(mask_full):
            return json.dumps({"error": "mask-not-found", "detail": f"mask not found: {mask_path}"})

        # Req 2.8: mask dimensions must match the original image
        try:
            from PIL import Image as _PIL_dim
            with _PIL_dim.open(full_path) as _img_orig:
                _orig_size = _img_orig.size
            with _PIL_dim.open(mask_full) as _img_mask:
                _mask_size = _img_mask.size
        except Exception as _e:
            return json.dumps({"error": "invalid-image", "detail": f"failed to read image dimensions: {str(_e)[:160]}"})
        if _orig_size != _mask_size:
            return json.dumps({
                "error": "mask-dimension-mismatch",
                "detail": f"mask {_mask_size[0]}x{_mask_size[1]} does not match image {_orig_size[0]}x{_orig_size[1]}"
            })

        with open(mask_full, "rb") as f:
            mask_b64 = base64.b64encode(f.read()).decode("ascii")

        for model_id in IMAGE_EDIT_MODELS:
            try:
                if model_id.startswith("stability."):
                    # Stable Image Inpaint / Search-and-Replace API.
                    # Body: {image: base64, mask: base64, prompt: str, output_format: png}
                    # search-and-replace는 mask 없이 prompt + search_prompt로 동작.
                    if "search-and-replace" in model_id:
                        body = {
                            "image": img_b64,
                            "prompt": prompt,
                            "search_prompt": prompt,  # 사용자가 별도 search_prompt 안 주면 prompt 재사용
                            "output_format": "png",
                        }
                    else:
                        body = {
                            "image": img_b64,
                            "mask": mask_b64,
                            "prompt": prompt,
                            "output_format": "png",
                        }
                elif model_id.startswith("amazon.titan"):
                    body = {
                        "taskType": "INPAINTING",
                        "inPaintingParams": {
                            "image": img_b64,
                            "maskImage": mask_b64,
                            "text": prompt,
                        },
                        "imageGenerationConfig": {"numberOfImages": 1, "quality": "standard"},
                    }
                else:  # nova-canvas
                    body = {
                        "taskType": "INPAINTING",
                        "inPaintingParams": {
                            "image": img_b64,
                            "maskImage": mask_b64,
                            "text": prompt,
                        },
                        "imageGenerationConfig": {"numberOfImages": 1},
                    }
                callable_id = _resolve_callable_model_id(model_id, aws_profile, bedrock_user)
                result = await gw.invoke_model(callable_id, body, timeout=60)
                if "error" in result:
                    last_error = f"{model_id}: {result['error'][:200]}"
                    continue
                images = result.get("images", [])
                if not images:
                    last_error = f"{model_id}: no images returned"
                    continue
                img_out = images[0] if isinstance(images[0], str) else (images[0].get("base64", "") if isinstance(images[0], dict) else "")
                if not img_out:
                    last_error = f"{model_id}: empty image data"
                    continue
                # Save
                # Always use a locally-existing directory for generated media.
                # Remote project_path may not exist locally; fall back to cwd.
                _local_root = _resolve_local_root(project_path)
                gen_dir = os.path.join(_local_root, ".generated")
                os.makedirs(gen_dir, exist_ok=True)
                ts = str(int(_t.time() * 1000))
                filename = f"inpaint-{ts}.png"
                output_path = os.path.join(gen_dir, filename)
                with open(output_path, "wb") as f:
                    f.write(base64.b64decode(img_out))
                try:
                    from PIL import Image as _PIL
                    with _PIL.open(output_path) as im:
                        aw, ah = im.size
                except Exception:
                    aw, ah = 0, 0
                return json.dumps({
                    "path": f".generated/{filename}",
                    "model": model_id,
                    "width": aw,
                    "height": ah,
                    "mode": "inpaint",
                })
            except Exception as e:
                last_error = f"{model_id}: {str(e)[:200]}"
                continue

        # [hint-image-route]
        _det = (last_error or "all inpaint models failed")[:200]
        _payload = {"error": "model-unavailable", "detail": _det}
        if any(t in _det for t in ("execute-api:Invoke", "principal identity", "HTTP 403", "HTTP 404")):
            _payload["hint"] = "현재 게이트웨이가 이미지 편집(inpaint) 라우트(/invoke-model)를 지원하지 않습니다. 관리자에게 활성화를 요청하세요."
        return json.dumps(_payload)

    # outpaint mode
    # Req 3.5: enforce 4096px max on the longer original edge
    try:
        from PIL import Image as _PIL_dim_op
        with _PIL_dim_op.open(full_path) as _src_img:
            _src_w, _src_h = _src_img.size
    except Exception as _e:
        return json.dumps({"error": "invalid-input", "detail": f"failed to read image dimensions: {str(_e)[:160]}"})
    if max(_src_w, _src_h) > 4096:
        return json.dumps({
            "error": "invalid-input",
            "detail": f"image dimension exceeds 4096px (got {_src_w}x{_src_h})",
        })

    # Req 3.7: direction is required, must be a list of 1..4 entries from the allowed set.
    # AGENT_TOOLS schema exposes "up"/"down" while spec uses "top"/"bottom"; accept both
    # and normalize to {left, right, top, bottom}.
    raw_direction = tool_input.get("direction", ["right"])
    if not isinstance(raw_direction, list) or not (1 <= len(raw_direction) <= 4):
        return json.dumps({
            "error": "invalid-parameter",
            "detail": "direction must be a list of 1-4 values (left/right/top/bottom)",
        })
    _dir_alias = {"up": "top", "down": "bottom", "top": "top", "bottom": "bottom",
                  "left": "left", "right": "right"}
    normalized = []
    seen = set()
    for d in raw_direction:
        if not isinstance(d, str) or d not in _dir_alias:
            return json.dumps({
                "error": "invalid-parameter",
                "detail": f"invalid direction value: {d!r} (allowed: left/right/top/bottom)",
            })
        nd = _dir_alias[d]
        if nd not in seen:
            seen.add(nd)
            normalized.append(nd)
    direction = normalized

    # Req 3.7: extend_pixels must be an integer in [1, 1024]
    raw_extend = tool_input.get("extend_pixels", 256)
    if isinstance(raw_extend, bool) or not isinstance(raw_extend, int):
        return json.dumps({
            "error": "invalid-parameter",
            "detail": "extend_pixels must be an integer in 1-1024",
        })
    extend_pixels = raw_extend
    if not (1 <= extend_pixels <= 1024):
        return json.dumps({"error": "invalid-parameter", "detail": "extend_pixels must be 1-1024"})

    for model_id in IMAGE_OUTPAINT_MODELS:
        try:
            if model_id.startswith("stability."):
                # Stable Image Outpaint API: {image, prompt, left/right/up/down, output_format}
                body = {
                    "image": img_b64,
                    "prompt": prompt,
                    "output_format": "png",
                }
                # direction을 Stable Outpaint 픽셀 인자로 변환
                for d in direction:
                    if d in ("left", "right", "up", "down", "top", "bottom"):
                        # Bedrock Stable Outpaint는 left/right/up/down — top→up, bottom→down 정규화
                        key = {"top": "up", "bottom": "down"}.get(d, d)
                        body[key] = int(extend_pixels)
            elif model_id.startswith("amazon.titan"):
                body = {
                    "taskType": "OUTPAINTING",
                    "outPaintingParams": {
                        "image": img_b64,
                        "text": prompt,
                        "outPaintingMode": "DEFAULT",
                    },
                    "imageGenerationConfig": {"numberOfImages": 1, "quality": "standard"},
                }
            else:
                body = {
                    "taskType": "OUTPAINTING",
                    "outPaintingParams": {
                        "image": img_b64,
                        "text": prompt,
                    },
                    "imageGenerationConfig": {"numberOfImages": 1},
                }
            callable_id = _resolve_callable_model_id(model_id, aws_profile, bedrock_user)
            result = await gw.invoke_model(callable_id, body, timeout=60)
            if "error" in result:
                last_error = f"{model_id}: {result['error'][:200]}"
                continue
            images = result.get("images", [])
            if not images:
                last_error = f"{model_id}: no images"
                continue
            img_out = images[0] if isinstance(images[0], str) else (images[0].get("base64", "") if isinstance(images[0], dict) else "")
            if not img_out:
                continue
            # Always use a locally-existing directory for generated media.
            # Remote project_path may not exist locally; fall back to cwd.
            _local_root = _resolve_local_root(project_path)
            gen_dir = os.path.join(_local_root, ".generated")
            os.makedirs(gen_dir, exist_ok=True)
            ts = str(int(_t.time() * 1000))
            filename = f"outpaint-{ts}.png"
            output_path = os.path.join(gen_dir, filename)
            with open(output_path, "wb") as f:
                f.write(base64.b64decode(img_out))
            try:
                from PIL import Image as _PIL
                with _PIL.open(output_path) as im:
                    aw, ah = im.size
            except Exception:
                aw, ah = 0, 0
            return json.dumps({
                "path": f".generated/{filename}",
                "model": model_id,
                "width": aw,
                "height": ah,
                "mode": "outpaint",
                "direction": direction,
            })
        except Exception as e:
            last_error = f"{model_id}: {str(e)[:200]}"
            continue

    # [hint-image-route]
    _det = (last_error or "all outpaint models failed")[:200]
    _payload = {"error": "model-unavailable", "detail": _det}
    if any(t in _det for t in ("execute-api:Invoke", "principal identity", "HTTP 403", "HTTP 404")):
        _payload["hint"] = "현재 게이트웨이가 이미지 편집(outpaint) 라우트(/invoke-model)를 지원하지 않습니다. 관리자에게 활성화를 요청하세요."
    return json.dumps(_payload)



def _execute_tool(tool_name: str, tool_input: dict, project_path: str = "", aws_profile: str = "", bedrock_user: str = "", template_id: str = "") -> str:
    """도구를 실행하고 결과를 문자열로 반환.

    template_id가 주어지고 tool_name이 generate_pptx면, 활성 템플릿을 해석해
    tool_input에 templatePath/templateId/styleProfile을 주입한다 (pptx-template-styling
    요구사항 5.1, 5.3). 이렇게 해야 모델이 도구를 *직접* 호출하는 단일 경로(run-agent)
    에서도 템플릿이 적용된다. (이전 버그: _force_generate_from_text 경로만 템플릿을
    처리하고, 모델이 generate_pptx를 직접 부르면 빈 Presentation()으로 생성됐음.)
    """
    # === Remote Bridge Routing ===
    _REMOTE_TOOLS = {"read_file", "write_file", "list_directory", "search_files", "run_command"}
    if not _BRIDGE_URL:
        _refresh_bridge_discovery()
    if tool_name in _REMOTE_TOOLS:
        _tgt_path = ""
        if isinstance(tool_input, dict):
            _tgt_path = str(tool_input.get("path") or tool_input.get("cwd") or "")
        def _looks_remote_path(_p):
            _p = _p or ""
            return (_p.startswith("/fsx/") or _p == "/fsx"
                    or _p.startswith("/home/") or "[SSH:" in _p)
        _path_remote = _looks_remote_path(_tgt_path)
        # 원격 경로거나 세션이 원격이면 브리지로 라우팅(플래그가 stale해도 원격 경로는 시도).
        if _BRIDGE_URL and (_path_remote or _bridge_is_remote()):
            _br = _call_bridge(tool_name, tool_input)
            if _br is not None:
                return _format_bridge_result(tool_name, _br)
            # 브리지 도달 불가 — 원격 경로면 로컬 폴백으로 인한 환각을 막기 위해 명확히 알린다.
            if _path_remote:
                return ("[원격 연결 필요] '" + _tgt_path[:120] + "' 경로는 원격(SSH) 서버에 있습니다. "
                        "현재 원격 세션을 읽을 수 없어 내용을 확인하지 못했습니다. 추측으로 내용을 지어내지 말고, "
                        "SSH 연결(원격 호스트)을 확인한 뒤 다시 시도하세요.")
        # 비원격 경로는 아래 로컬 처리로 폴백(기존 동작).

    # === 템플릿 주입 (요구사항 5.1, 5.3) — generate_pptx 직접 호출 경로 ===
    # 모델이 도구를 직접 호출하면 tool_input엔 모델이 만든 {title, slides}만 있고
    # templatePath가 없다. 활성 템플릿을 여기서 해석해 주입한다. 무템플릿이면
    # tool_input 불변(기존 동작 보존, 요구사항 5.2).
    if tool_name == "generate_pptx" and template_id and isinstance(tool_input, dict) \
            and not tool_input.get("templatePath"):
        try:
            _store_root = None
            _tm = _load_template_manager()
            if _tm is not None:
                _store_root = _tm.resolve_template_store_root()
            _tpl_path, _tpl_profile, _tpl_used = _resolve_active_template(template_id, _store_root)
            if _tpl_used:
                tool_input = dict(tool_input)  # 원본 mutate 방지
                tool_input["templatePath"] = _tpl_path
                tool_input["templateId"] = template_id
                if _tpl_profile:
                    tool_input["styleProfile"] = _tpl_profile
                print(f"[ExecuteTool] generate_pptx에 템플릿 주입 — templateId={str(template_id)[:80]}")
        except Exception as _te:
            # 템플릿 해석 실패는 무템플릿으로 격리 (요구사항 9)
            print(f"[ExecuteTool] 템플릿 주입 실패 → 무템플릿 진행: {str(_te)[:200]}")

    # Async media generation tools
    if tool_name in ("generate_image", "generate_pdf", "generate_pptx", "generate_xlsx", "generate_docx", "edit_image", "generate_native_diagram"):
        try:
            import asyncio as _asyncio
            if tool_name == "generate_image":
                return _asyncio.run(_tool_generate_image(tool_input, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user))
            if tool_name == "generate_pdf":
                return _asyncio.run(_tool_generate_pdf(tool_input, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user))
            if tool_name == "generate_pptx":
                return _asyncio.run(_tool_generate_pptx(tool_input, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user))
            if tool_name == "generate_xlsx":
                return _asyncio.run(_tool_generate_xlsx(tool_input, project_path))
            if tool_name == "generate_docx":
                return _asyncio.run(_tool_generate_docx(tool_input, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user))
            if tool_name == "edit_image":
                return _asyncio.run(_tool_edit_image(tool_input, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user))
            if tool_name == "generate_native_diagram":
                # 독립 실행 진입점. 헬퍼는 개별 kwargs(diagram_type/title/content)를 받고
                # 다른 generate_* 도구와 동일하게 {path: ".generated/...", model, width,
                # height, sizeBytes} JSON을 반환한다 → GatewayToolNode의 verified_files
                # 실측(path 추출)이 그대로 동작한다 (요구사항 1.6/3.7).
                return _asyncio.run(_tool_generate_native_diagram(
                    diagram_type=str(tool_input.get("diagram_type", "tree")),
                    title=str(tool_input.get("title", "Diagram")),
                    content=str(tool_input.get("content", "")),
                    project_path=project_path,
                ))
        except Exception as e:
            return json.dumps({"error": "tool-execution-failed", "detail": str(e)[:300]})


    try:
        if tool_name == "read_file":
            path = tool_input["path"]
            if not os.path.isabs(path) and project_path:
                path = os.path.join(project_path, path)
            if not os.path.exists(path):
                return f"파일 없음: {path}"
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            _rf_max = int(os.environ.get("AE_READ_FILE_MAX", "120000"))
            if len(content) > _rf_max:
                content = content[:_rf_max] + f"\n... (총 {len(content)}자, {_rf_max}자까지 표시)"
            return content

        elif tool_name == "write_file":
            path = tool_input["path"]
            if not os.path.isabs(path) and project_path:
                path = os.path.join(project_path, path)
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, 'w', encoding='utf-8') as f:
                f.write(tool_input["content"])
            return f"파일 저장 완료: {path} ({len(tool_input['content'])}자)"

        elif tool_name == "list_directory":
            path = tool_input["path"]
            if not os.path.isabs(path) and project_path:
                path = os.path.join(project_path, path)
            if not os.path.isdir(path):
                return f"디렉토리 없음: {path}"
            entries = os.listdir(path)
            result = []
            for e in sorted(entries):
                if e.startswith('.') and e not in ('.env', '.gitignore'):
                    continue
                fp = os.path.join(path, e)
                kind = "DIR" if os.path.isdir(fp) else "FILE"
                result.append(f"  {kind}  {e}")
            return f"{path}/ ({len(result)}개)\n" + "\n".join(result[:100])

        elif tool_name == "run_command":
            cmd = tool_input["command"]
            cwd = tool_input.get("cwd", project_path or os.getcwd())
            result = subprocess.run(
                cmd, shell=True, capture_output=True, text=True,
                timeout=30, cwd=cwd,
                env={**os.environ, "PATH": os.environ.get("PATH", "")},
            )
            output = result.stdout + result.stderr
            _rc_max = int(os.environ.get("AE_RUN_CMD_MAX", "40000"))
            if len(output) > _rc_max:
                output = output[:_rc_max] + f"\n... (출력 잘림, 총 {len(output)}자 중 {_rc_max}자 표시)"
            return output or "(출력 없음)"

        elif tool_name == "search_files":
            query = tool_input["query"]
            path = tool_input["path"]
            if not os.path.isabs(path) and project_path:
                path = os.path.join(project_path, path)
            pattern = tool_input.get("file_pattern", "")
            include = f"--include='{pattern}'" if pattern else ""
            cmd = f"grep -rn {include} --color=never '{query}' '{path}' 2>/dev/null | head -50"
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=10)
            return result.stdout or "검색 결과 없음"

        else:
            return f"알 수 없는 도구: {tool_name}"
    except subprocess.TimeoutExpired:
        return "명령 실행 시간 초과 (30초)"
    except Exception as e:
        return f"도구 실행 오류: {str(e)}"


# GatewayClient 캐시 — 동일 profile+user 조합은 재사용
_gw_cache = {}


def _is_code_related(prompt: str) -> bool:
    """프롬프트가 코드/프로젝트 관련인지 판단."""
    p = prompt.lower().strip()
    # 코드 관련 키워드
    code_keywords = [
        'code', 'function', 'class', 'import', 'error', 'bug', 'fix',
        'implement', 'refactor', 'test', 'deploy', 'build', 'compile',
        '코드', '함수', '클래스', '에러', '버그', '수정', '구현', '리팩토링',
        '파일', '모듈', '컴포넌트', '테스트', '배포', '빌드', '변수',
        'api', 'endpoint', 'database', 'query', 'schema', 'migration',
        'this file', 'this project', '이 파일', '이 프로젝트', '현재',
        '경로', '폴더', '디렉토리', '열린', '오픈', 'path', 'directory',
        '에디터', 'editor', 'project', '프로젝트',
        '.js', '.py', '.ts', '.css', '.html', '.json',
    ]
    for kw in code_keywords:
        if kw in p:
            return True
    # 200자 이상이면 코드 관련일 가능성 높음
    if len(p) > 200:
        return True
    return False


def _build_messages(chat_history: list, current_prompt: str, session_id: str = "") -> list:
    """ConversationMemory를 통해 messages 구성."""
    from ai_engine.rag.conversation_memory import get_memory
    mem = get_memory()
    messages, _ = mem.build_messages(session_id or "default", chat_history, current_prompt)
    return messages

async def _maybe_summarize(session_id: str, chat_history: list, gw):
    """대화가 길어지면 비동기로 요약 체크포인트 생성."""
    try:
        from ai_engine.rag.conversation_memory import get_memory
        mem = get_memory()
        _, needs = mem.build_messages(session_id, chat_history, "")
        if needs:
            await mem.summarize_and_checkpoint(session_id, chat_history, gw)
    except Exception as e:
        print(f"[Memory] 요약 트리거 실패: {e}")


def _get_gw(aws_profile, bedrock_user):
    key = f"{aws_profile}:{bedrock_user}"
    if key not in _gw_cache:
        from ai_engine.gateway_module import GatewayClient
        _gw_cache[key] = GatewayClient(
            gateway_url=os.environ.get("GATEWAY_URL", "https://5l764dh7y9.execute-api.us-west-2.amazonaws.com/v1"),
            aws_profile=aws_profile,
            region=os.environ.get("AWS_REGION", "us-west-2"),
            bedrock_user=bedrock_user,
        )
    gw = _gw_cache[key]
    # 주입된 자격증명이 있으면 캐시 만료하지 않음
    if not hasattr(gw, '_injected_creds') or not gw._injected_creds:
        gw._cred_time = 0
    return gw


# ─── Callable Model ID Resolution ──────────────────────────────────────
# Bedrock 모델 ID는 3가지 형태가 있음:
#  1) ON_DEMAND: 'nvidia.nemotron-nano-12b-v2' (prefix 없이 직접 호출)
#  2) INFERENCE_PROFILE (CRIS): 'us.anthropic.claude-opus-4-7' (us./eu./global. prefix 필수)
#  3) 둘 다 지원: 이 경우 CRIS 우선 (성능·리전 분산)
#
# foundation model 목록의 inferenceTypesSupported를 보고 올바른 형태를 선택한다.
# 캐시는 profile 부팅 후 list_foundation_models로 한 번만 조회 (프로세스 생애).

_model_inference_type_cache = {}  # key: "{profile}:{user}", val: {modelId: [types]}


def _load_inference_types(aws_profile, bedrock_user):
    key = f"{aws_profile}:{bedrock_user}"
    if key in _model_inference_type_cache:
        return _model_inference_type_cache[key]
    try:
        import boto3
        session = boto3.Session(profile_name=aws_profile, region_name=os.environ.get("AWS_REGION", "us-west-2"))
        client = session.client("bedrock")
        resp = client.list_foundation_models()
        cache = {}
        for m in resp.get("modelSummaries", []):
            mid = m.get("modelId", "")
            cache[mid] = m.get("inferenceTypesSupported", []) or []
        _model_inference_type_cache[key] = cache
        return cache
    except Exception as e:
        print(f"[ModelResolver] list_foundation_models 실패: {e}")
        _model_inference_type_cache[key] = {}
        return {}


# ─── Specialized Model Routing Matrix (모듈 스코프) ───────────────────
# 사용자의 "Option A": task type별로 자동 최적 모델 선택, 사용자 선택은 항상 우선.
# 이 리스트들은 게이트웨이가 활성화한 Claude 모델 우선순위 (newest-first).
# 게이트웨이가 일부 모델만 활성화한 경우, _resolve_callable_model_id가 호출 시점에
# us./eu. prefix 처리를 담당한다.
_LATEST_OPUS_IDS = [
    "anthropic.claude-opus-4-7-20251015-v1:0",
    "anthropic.claude-opus-4-20250514-v1:0",
    "anthropic.claude-3-opus-20240229-v1:0",
    # opus 권한이 없는 계정(라이브 실측: opus 전부 model_denied)에서도 동작하도록
    # 활성 sonnet-4-5를 최종 폴백으로 둔다. opus 허용 계정은 위 opus를 그대로 사용.
    "anthropic.claude-sonnet-4-5-20250929-v1:0",
]
_LATEST_SONNET_IDS = [
    # sonnet-4-5: 라이브 게이트웨이 실측 활성(2026-07). 4-6은 일부 계정에서
    # model_denied 되므로 활성 확실한 4-5를 최우선 폴백으로 둔다.
    "anthropic.claude-sonnet-4-5-20250929-v1:0",
    "anthropic.claude-sonnet-4-6-20250929-v1:0",
    "anthropic.claude-sonnet-4-20250514-v1:0",
    "anthropic.claude-3-5-sonnet-20241022-v2:0",
]
_LATEST_HAIKU_IDS = [
    "anthropic.claude-haiku-4-5-20251001-v1:0",
    "anthropic.claude-3-5-haiku-20241022-v1:0",
]

# === 동적 게이트웨이 카탈로그 캐시 ===
# /api/models 응답을 통해 게이트웨이가 실제 활성화한 모델 ID 전체를 메모리에 보관.
# _specialized_model_for_task가 매트릭스 고정 풀 외에도 이 캐시에서 task에 맞는
# provider를 찾아 활용. Gemma처럼 Anthropic이 아닌 새 모델이 게이트웨이에 추가되면
# 자동으로 라우팅 후보가 됨 (코드 변경 없이 매트릭스 확장).
_GATEWAY_MODEL_CACHE = {
    "models": [],     # [{id, provider, name, capabilities}, ...]
    "last_fetched": 0,
    "ttl": 300,       # 5분
}

# === 라우팅 denylist ===
# 게이트웨이가 `model_denied`("not in allowed list")로 거부한 모델을 기억해 두고
# 자동 라우팅 후보에서 제외한다. 이렇게 해야 Gemma `-pt-`(사전학습 base, chat 미허용)
# 같은 미허용 모델로 라우팅돼 단계 전체가 실패하는 일을 막는다.
# - 정적 시드: instruction-tuned가 아닌 base 변형(`-pt-`)은 어떤 게이트웨이에서도
#   chat 호출이 거부되므로 사전 차단.
# - 런타임 학습: 스트림 에러에서 `model_denied`를 감지하면 해당 id를 추가.
_DENIED_MODEL_IDS = set()           # 런타임에 관찰된 정확한 거부 id
_DENIED_MODEL_PATTERNS = ["-pt-v", "-pt:"]  # base(pretrained) 변형 — chat 미허용


def _normalize_model_key(model_id: str) -> str:
    """prefix(us./eu./global.)를 떼어 비교용 raw id 반환."""
    if not model_id:
        return ""
    mid = model_id
    for prefix in ("us.", "eu.", "global."):
        if mid.startswith(prefix):
            mid = mid[len(prefix):]
            break
    return mid.lower()


def _model_is_denied(model_id: str) -> bool:
    """모델이 denylist(정확 id) 또는 base-변형 패턴에 해당하면 True."""
    if not model_id:
        return False
    raw = _normalize_model_key(model_id)
    if raw in _DENIED_MODEL_IDS:
        return True
    for pat in _DENIED_MODEL_PATTERNS:
        if pat in raw:
            return True
    return False


def _record_denied_model(model_id: str) -> None:
    """게이트웨이가 거부한 모델 id를 denylist에 등록(런타임 학습)."""
    raw = _normalize_model_key(model_id)
    if raw and raw not in _DENIED_MODEL_IDS:
        _DENIED_MODEL_IDS.add(raw)
        print(f"[ModelRouter] denylist 추가 — {raw} (게이트웨이 model_denied)")


def _extract_denied_model_from_error(detail: str) -> str:
    """`model_denied` 에러 메시지에서 거부된 모델 id를 추출. 없으면 ''."""
    if not detail:
        return ""
    import re as _re
    # 예: "model us.google.gemma-3-27b-pt-v1:0 not in allowed list"
    m = _re.search(r"model\s+([A-Za-z0-9._:\-]+)\s+not in allowed", detail)
    if m:
        return m.group(1)
    return ""


def _maybe_record_denied_from_error(detail: str) -> None:
    """스트림/호출 에러 문자열에 model_denied가 있으면 해당 모델을 denylist에 추가."""
    if not detail:
        return
    low = detail.lower()
    if "model_denied" in low or "not in allowed list" in low:
        mid = _extract_denied_model_from_error(detail)
        if mid:
            _record_denied_model(mid)


def _update_gateway_model_cache(catalog: dict):
    """list_models 응답을 받아 캐시 갱신. catalog는 {provider: [{id, name, capabilities}]} 형태."""
    import time as _t
    flat = []
    for prov, items in (catalog or {}).items():
        for m in items or []:
            mid = m.get("id") or ""
            if not mid:
                continue
            flat.append({
                "id": mid,
                "provider": prov,
                "name": m.get("name", mid),
                "capabilities": m.get("capabilities", {}),
            })
    _GATEWAY_MODEL_CACHE["models"] = flat
    _GATEWAY_MODEL_CACHE["last_fetched"] = _t.time()
    print(f"[ModelCache] {len(flat)} models cached from gateway catalog")


# ─────────────────────────────────────────────────────────────────
# OpenAI 모델 라우팅 (gateway-openai-models, 요구사항 5)
# 선택된 모델이 OpenAI provider면 /openai/responses(동기)→/openai/responses-jobs(비동기)
# 로 라우팅하고, 응답을 OpenAI_Response_Adapter로 Converse 형식으로 변환한다.
# 비-OpenAI는 기존 Bedrock 경로(gw.converse/stream)로 그대로(바이트 동일).
# ─────────────────────────────────────────────────────────────────
def _openai_model_ids() -> set:
    """현재 카탈로그 캐시에서 provider=="OpenAI"인 모델 id 집합."""
    ids = set()
    for m in _GATEWAY_MODEL_CACHE.get("models", []):
        if (m.get("provider") or "") == "OpenAI" and m.get("id"):
            ids.add(m["id"])
    return ids


def is_openai_model(model_id: str, openai_ids: set | None = None) -> bool:
    """카탈로그 멤버십 우선, 보조로 'openai.' prefix로 OpenAI 모델 판정."""
    if not model_id:
        return False
    if openai_ids is None:
        openai_ids = _openai_model_ids()
    return model_id in openai_ids or model_id.startswith("openai.")


async def route_openai_chat(gw, model_id, messages, system_prompt="", timeout=120):
    """OpenAI 모델 채팅 — 동기 우선, 타임아웃/실패 시 비동기 폴백 → Converse 변환.

    반환: Bedrock Converse 형식 dict (openai_adapter.to_converse 결과).
    동기+비동기 모두 실패 시 예외를 그대로 전파(부분 응답 미전달, 요구사항 7.4).
    """
    try:
        from ai_engine import openai_adapter as _ad
        from ai_engine.gateway_module import SyncTimeout as _SyncTimeout
    except ImportError:
        import openai_adapter as _ad
        from gateway_module import SyncTimeout as _SyncTimeout
    try:
        raw = await gw.openai_responses_sync(
            model_id, messages, system_prompt=system_prompt, timeout=timeout)
    except _SyncTimeout:
        print(f"[OpenAIRoute] 동기 타임아웃 → 비동기 잡 폴백: {model_id}")
        raw = await gw.openai_responses_job_submit_and_poll(
            model_id, messages, system_prompt=system_prompt,
            poll_interval=5, max_wait=300)
    return _ad.to_converse(raw)


async def route_openai_agent(gw, model_id, messages, system_prompt="", project_path="",
                             aws_profile="", bedrock_user="", template_id="",
                             max_iters=8, timeout=120):
    """OpenAI 함수호출(tool) 실행 루프 — GPT가 스스로 도구를 호출해 자율 동작.

    OpenAI Responses 라우트에 tools를 실어 보내고, 모델이 반환한 function_call을
    _execute_tool로 실제 실행한 뒤 결과를 input에 누적해 다시 호출한다. 도구 호출이
    없으면 최종 텍스트를 반환한다. 이렇게 해야 GPT 5.5/5.4도 단일 모드에서 검색→
    읽기→추론→생성까지 자율 수행한다(Claude의 tool loop와 동등).

    반환: {"text": <최종 텍스트>, "verified_files": [{path, absPath, size, tool}]}
    어떤 예외도 호출자(엔드포인트)가 잡아 force-generate 폴백으로 복구한다.
    """
    try:
        from ai_engine import openai_adapter as _ad
    except ImportError:
        import openai_adapter as _ad

    tools = _agent_tools_to_openai()
    input_items = gw._to_openai_input(messages)
    # 도구 사용 지침 — 흐름도/구조도는 네이티브 편집가능 다이어그램으로(이미지 생성 불요).
    # 이미지 게이트웨이가 불안정한 환경에서도 시각자료가 확실히 생성되도록 유도한다.
    _tool_guide = (
        "\n\n[도구 사용 지침 — 매우 중요]\n"
        "1. 파일 생성을 요청받으면 반드시 generate_pptx/generate_pdf/generate_docx/generate_xlsx "
        "도구를 실제로 호출하세요. 텍스트로만 '생성했다'고 답하지 마세요.\n"
        "2. PPTX/PDF는 시각적으로 풍부해야 합니다. 표지 슬라이드와 주요 섹션 슬라이드에는 "
        "반드시 imagePrompt 필드로 고품질 이미지를 요청하세요(영어로 상세히, 예: "
        "'modern flat illustration of cloud architecture, blue palette, professional'). "
        "이미지는 시스템이 최고 품질 모델로 자동 생성합니다.\n"
        "3. 흐름도/구조도/아키텍처/순서도 같은 '다이어그램'은 imagePrompt 대신 해당 슬라이드의 "
        "bullets에 단계/구성요소를 순서대로 적으세요. 시스템이 PowerPoint 네이티브 편집가능 "
        "도형(박스+화살표)으로 자동 변환합니다. (이미지 슬라이드와 다이어그램 슬라이드를 함께 구성)\n"
        "4. 필요하면 read_file/list_directory/search_files로 프로젝트를 먼저 파악한 뒤 작성하세요."
    )
    system_prompt = (system_prompt or "") + _tool_guide
    loop = asyncio.get_event_loop()
    text_acc = ""
    verified_files = []
    seen_abs = set()

    def _track_tool_result(tool_name, result_str):
        # 생성 도구 결과 JSON에서 path/absPath를 뽑아 디스크 검증 후 기록.
        try:
            obj = json.loads(result_str)
        except (ValueError, TypeError):
            return
        if not isinstance(obj, dict):
            return
        rel = obj.get("path") or ""
        _abs = obj.get("absPath") or (_resolve_relative_for_verify(rel, project_path) if rel else "")
        if not _abs:
            return
        try:
            if os.path.isfile(_abs) and os.path.getsize(_abs) > 0 and _abs not in seen_abs:
                seen_abs.add(_abs)
                verified_files.append({
                    "path": rel or os.path.basename(_abs),
                    "absPath": _abs,
                    "size": os.path.getsize(_abs),
                    "tool": tool_name,
                })
        except OSError:
            pass

    for _iter in range(max_iters):
        body = {"model": model_id, "input": input_items, "tools": tools, "tool_choice": "auto"}
        if system_prompt:
            body["instructions"] = system_prompt
        raw = await gw.openai_responses_call(body, timeout=timeout)
        _t = _ad.extract_text(raw)
        if _t:
            text_acc += (("\n" + _t) if text_acc else _t)
        calls = _ad.extract_function_calls(raw)
        if not calls:
            break
        for c in calls:
            try:
                args = json.loads(c["arguments"]) if isinstance(c["arguments"], str) else (c["arguments"] or {})
            except (ValueError, TypeError):
                args = {}
            if not isinstance(args, dict):
                args = {}
            try:
                result_str = await loop.run_in_executor(
                    None, _execute_tool, c["name"], args, project_path,
                    aws_profile, bedrock_user, template_id,
                )
            except Exception as _ee:
                result_str = json.dumps({"error": "tool-failed", "detail": str(_ee)[:200]})
            _track_tool_result(c["name"], result_str)
            # OpenAI Responses 멀티턴 — function_call + function_call_output을 input에 누적.
            input_items.append({
                "type": "function_call",
                "call_id": c["call_id"],
                "name": c["name"],
                "arguments": c["arguments"] if isinstance(c["arguments"], str) else json.dumps(args),
            })
            input_items.append({
                "type": "function_call_output",
                "call_id": c["call_id"],
                "output": str(result_str)[:6000],
            })
    return {"text": text_acc, "verified_files": verified_files}


def _gateway_models_by_keyword(*keywords) -> list:
    """캐시에서 keyword가 model id에 포함된 것들 반환 (대소문자 무시)."""
    if not _GATEWAY_MODEL_CACHE["models"]:
        return []
    results = []
    for m in _GATEWAY_MODEL_CACHE["models"]:
        mid_low = m["id"].lower()
        for kw in keywords:
            if kw.lower() in mid_low:
                results.append(m["id"])
                break
    return results

# === 비-Claude 모델 풀 ===
# 게이트웨이가 활성화한 다양한 모델을 task별로 최적 활용. 도구 호출이 필요없는
# 작업(요약/분류/단순 채팅)에는 비용 효율적 모델을, 추론/수학/긴 컨텍스트에는
# DeepSeek/Qwen 같은 특화 모델을 시도. 호출 실패 시 Claude로 fallback.
_REASONING_IDS = [
    "deepseek.r1-v1:0",
    "deepseek.v3-v1:0",
    "qwen.qwen3-235b-a22b-2507-v1:0",
] + _LATEST_OPUS_IDS  # fallback chain

_LIGHT_CHAT_IDS = [
    # Gemma 3 (Google) — 사용자 게이트웨이 카탈로그에 노출됨. Bedrock에서
    # 호출 가능한 정확한 ID는 게이트웨이가 활성화한 형태에 따라 다를 수 있어
    # 여러 후보를 fallback chain으로 둠 (앞의 것 실패하면 자동으로 다음 시도).
    "google.gemma-3-27b-pt-v1:0",
    "google.gemma-3-12b-it-v1:0",
    "google.gemma-3-4b-it-v1:0",
    # Gemma 2 — 구버전이지만 일부 게이트웨이는 이것만 활성
    "google.gemma-2-27b-it-v1:0",
    # Amazon Nova 경량
    "amazon.nova-lite-v1:0",
    "amazon.nova-micro-v1:0",
    # Meta / Mistral 경량
    "meta.llama3-3-70b-instruct-v1:0",
    "mistral.mistral-small-2402-v1:0",
] + _LATEST_HAIKU_IDS

_LONG_CONTEXT_IDS = [
    "qwen.qwen3-235b-a22b-2507-v1:0",
    "deepseek.v3-v1:0",
] + _LATEST_SONNET_IDS

_KOREAN_FRIENDLY_IDS = [
    "qwen.qwen3-235b-a22b-2507-v1:0",
] + _LATEST_SONNET_IDS


def _module_is_tool_capable(model_id: str) -> bool:
    """chat 모델이 Bedrock toolConfig를 안정적으로 지원하는지 확인 (모듈 스코프).

    - Claude 전체 → 도구 호출 안정
    - Mistral-Large / Pixtral → 도구 호출 안정
    - Nova-Pro → 도구 호출 가능 (Nova Lite/Micro는 제한적)
    - Llama / DeepSeek / Cohere / 기타 → 도구 호출 미지원/불안정
    """
    if not model_id:
        return False
    mid = model_id.lower()
    if "claude" in mid:
        return True
    # OpenAI(GPT 5.x 등) — 게이트웨이 OpenAI Responses 라우트 + 자체 함수호출 실행
    # 루프(route_openai_agent)로 도구 호출을 지원하므로 사용자 선택을 보존한다.
    if mid.startswith("openai.") or mid.startswith("us.openai.") or "gpt-" in mid:
        return True
    if "mistral-large" in mid or "pixtral" in mid:
        return True
    if "nova-pro" in mid:
        return True
    return False


def _module_is_vision_capable(model_id: str) -> bool:
    """모델이 이미지 입력(비전)을 지원하는지 확인."""
    if not model_id:
        return False
    mid = model_id.lower()
    # Claude 3.5 이상은 비전 지원, Haiku 4.5도 비전 지원
    if "claude" in mid:
        return True
    # OpenAI GPT 5.x는 멀티모달(비전) 지원
    if mid.startswith("openai.") or mid.startswith("us.openai.") or "gpt-" in mid:
        return True
    # Pixtral / Nova-Pro는 비전 지원
    if "pixtral" in mid or "nova-pro" in mid:
        return True
    return False


def _module_with_prefix(mid: str) -> str:
    """단순 us./eu. prefix 부여 — aws_profile/bedrock_user 없이 호출될 때 사용."""
    if not mid:
        return mid
    if mid.startswith("us.") or mid.startswith("eu.") or mid.startswith("global."):
        return mid
    return f"us.{mid}"


# task_type → 후보 ID 리스트 매핑.
# 각 리스트의 첫 번째가 최선, 이후는 fallback (newest-first).
# Claude는 도구 호출/장문 추론에서 안정적이므로 도구 필요 작업에는 우선.
# 비-Claude(DeepSeek/Qwen/Nova/Llama)는 도구 호출 부담 없는 작업에 활용.
_TASK_MODEL_TIERS = {
    # 도구 호출 + 코드 추론 — Claude Opus 우선
    "code":          _LATEST_OPUS_IDS,
    "code_analysis": _LATEST_OPUS_IDS,
    "refactor":      _LATEST_OPUS_IDS,
    "planner":       _LATEST_OPUS_IDS,
    "merger":        _LATEST_OPUS_IDS,
    # 파일 생성 + 콘텐츠 보강 — Sonnet (도구 호출 + 빠름)
    "file_generation": _LATEST_SONNET_IDS,
    "enrich":        _LATEST_SONNET_IDS,
    "vision_input":  _LATEST_SONNET_IDS,   # Claude Sonnet은 비전 지원
    # 단순 작업 — Haiku 우선이지만 비-Claude 경량 모델도 시도
    "simple_qa":     _LIGHT_CHAT_IDS,
    "general_chat":  _LIGHT_CHAT_IDS,
    "intent_classifier": _LATEST_HAIKU_IDS,  # 분류는 Haiku 안정적 (낮은 latency)
    "summarize":     _LIGHT_CHAT_IDS,
    "handoff":       _LIGHT_CHAT_IDS,
    # 순수 추론/수학/논리 (도구 X) — DeepSeek-R1 우선, Qwen, Opus fallback
    "reasoning":     _REASONING_IDS,
    "math":          _REASONING_IDS,
    "logic":         _REASONING_IDS,
    # 긴 문서 처리 — Qwen3-235B (대용량 컨텍스트), DeepSeek-V3
    "long_context":  _LONG_CONTEXT_IDS,
    "doc_analysis":  _LONG_CONTEXT_IDS,
    # 한국어 친화 — Qwen3 (다국어 강함), Sonnet
    "korean":        _KOREAN_FRIENDLY_IDS,
    "translation":   _KOREAN_FRIENDLY_IDS,
}


def _specialized_model_for_task(
    task_type: str,
    user_pref: str = "",
    capabilities_hint=None,
    aws_profile: str = "",
    bedrock_user: str = "",
) -> str:
    """task_type 기반 자동 모델 라우팅 — 사용자 override는 항상 최우선.

    우선순위:
    1. user_pref가 설정되어 있고, 도구 호출 가능하며, capabilities_hint(예: vision)
       조건을 만족하면 → user_pref 그대로 (사용자 선택 항상 승리).
    2. _TASK_MODEL_TIERS의 task별 후보 리스트에서 첫 번째 ID 선택.
    3. 매칭 실패 시 user_pref를 최종 fallback으로 반환.

    Args:
        task_type: code | code_analysis | refactor | file_generation | simple_qa |
                   general_chat | vision_input | planner | merger | enrich |
                   intent_classifier | summarize | handoff
        user_pref: 사용자가 명시적으로 선택한 모델 ID (raw 또는 prefix 포함)
        capabilities_hint: 추가 요구사항. 예: ["vision"] — 비전 입력이 있을 때
        aws_profile, bedrock_user: 있으면 _resolve_callable_model_id로 정확한 prefix.
                                    없으면 _module_with_prefix로 us. 단순 부여.

    Returns:
        호출 가능한 모델 ID (prefix 적용됨).
    """
    caps = set(capabilities_hint or [])
    needs_vision = "vision" in caps

    def _finalize(mid: str) -> str:
        if not mid:
            return mid
        if aws_profile or bedrock_user:
            try:
                return _resolve_callable_model_id(mid, aws_profile, bedrock_user)
            except Exception:
                return _module_with_prefix(mid)
        return _module_with_prefix(mid)

    # 1) 사용자 선택이 모든 조건을 만족하면 그대로 사용 (Option A 핵심)
    if user_pref:
        tool_ok = _module_is_tool_capable(user_pref)
        vision_ok = (not needs_vision) or _module_is_vision_capable(user_pref)
        if tool_ok and vision_ok:
            picked = _finalize(user_pref)
            print(f"[ModelRouter] task={task_type} user={user_pref} → picked={picked} (user override)")
            return picked

    # 2) task 매트릭스 lookup
    candidates = _TASK_MODEL_TIERS.get((task_type or "").lower())
    # 매트릭스 후보 + 동적 카탈로그에서 task별 keyword 매칭으로 발견된 모델 합치기.
    # 이렇게 하면 게이트웨이에 새 모델(Gemma, 새 Llama 등)이 추가되면 자동으로 후보로
    # 들어옴. 기존 정적 풀이 우선이고 동적 발견은 fallback 위치.
    task_lower = (task_type or "").lower()
    dynamic_extra = []
    if task_lower in ("simple_qa", "general_chat", "summarize", "handoff"):
        # 경량 채팅 — Gemma 3 (멀티모달·경량 추론), Llama-small, Phi, Nova-Lite/Micro 등
        dynamic_extra = _gateway_models_by_keyword(
            "gemma-3", "gemma-2", "gemma", "phi", "ministral",
            "llama-3-2-3b", "llama-3-1-8b", "nova-lite", "nova-micro",
        )
    elif task_lower in ("reasoning", "math", "logic"):
        dynamic_extra = _gateway_models_by_keyword(
            "deepseek", "r1", "qwen3", "o1", "openai",
        )
    elif task_lower in ("translation", "korean", "long_context", "doc_analysis"):
        dynamic_extra = _gateway_models_by_keyword(
            "qwen3", "qwen-2", "command-r-plus", "mistral-large",
        )
    elif task_lower in ("code", "code_analysis", "refactor"):
        # 코드는 도구 호출 안정성이 핵심 — Claude 우선이지만 Codestral도 후보
        dynamic_extra = _gateway_models_by_keyword(
            "codestral", "deepseek-coder", "qwen-coder",
        )

    if candidates:
        # 동적 발견 모델은 정적 후보 다음 fallback으로 (정적이 우선이라 안정성 유지)
        merged = list(candidates)
        for d in dynamic_extra:
            if d not in merged:
                merged.append(d)
        candidates = merged
        # vision 필요 시 vision-capable한 후보만 (Claude/Pixtral/Nova-Pro)
        if needs_vision:
            candidates = [c for c in candidates if _module_is_vision_capable(c)] or candidates
        # 도구 호출이 필요한 task인데 후보 중 비-도구 모델이 섞여있으면 도구 가능한 것만.
        _tool_required_tasks = {
            "planner", "merger", "code", "code_analysis", "refactor",
            "file_generation", "enrich",
        }
        if (task_type or "").lower() in _tool_required_tasks:
            tool_candidates = [c for c in candidates if _module_is_tool_capable(c)]
            if tool_candidates:
                candidates = tool_candidates
        # 게이트웨이가 거부(model_denied)했거나 chat 미허용 base 변형(`-pt-`)인 후보 제외.
        # 전부 걸러지면 Claude로 안전 폴백 — 단계 실패 대신 항상 호출 가능한 모델 보장.
        _allowed = [c for c in candidates if not _model_is_denied(c)]
        if _allowed:
            candidates = _allowed
        else:
            candidates = list(_LATEST_SONNET_IDS)
        chosen = candidates[0]
        picked = _finalize(chosen)
        print(f"[ModelRouter] task={task_type} user={user_pref or '∅'} → picked={picked} (matrix; dyn_extra={len(dynamic_extra)})")
        return picked

    # 3) 매트릭스에도 없고 user_pref도 없으면 Sonnet 최종 fallback
    fallback = _finalize(user_pref or _LATEST_SONNET_IDS[0])
    print(f"[ModelRouter] task={task_type} user={user_pref or '∅'} → picked={fallback} (final fallback)")
    return fallback


def _finalize_route_to_claude(current_model: str = "", aws_profile: str = "", bedrock_user: str = "") -> str:
    """미허용(model_denied)/불안정 모델을 항상 호출 가능한 Claude로 재라우팅.

    file_generation task 풀(Sonnet — 도구 호출 가능 + 빠름)에서 선택한다. 현재 모델이
    이미 Claude면 그대로 둔다(호출부에서 동일 여부 비교). denylist를 반영하므로 거부된
    Sonnet 변형이 있으면 다음 후보로 자동 폴백된다.
    """
    if current_model and "claude" in current_model.lower() and not _model_is_denied(current_model):
        return current_model
    return _specialized_model_for_task(
        "file_generation", "", aws_profile=aws_profile, bedrock_user=bedrock_user,
    )


# === SSE 스트림 heartbeat 래퍼 ===
# 모델이 오래 "thinking" 중이거나 도구가 길게 실행되면 게이트웨이 스트림에서
# 수십 초~수 분간 이벤트가 안 나올 수 있다. 클라이언트의 idle 워치독(기본 180초)이
# 이를 "끊김"으로 오판해 작업을 중단시키는 문제를 막기 위해, idle 구간마다
# {"type":"heartbeat"} 합성 이벤트를 주입해 스트림을 살아있게 유지한다.
# (사용자 요구: 오래 걸려도 thinking 중이면 진행 표시하며 계속)
_SSE_HEARTBEAT_SECONDS = float(os.environ.get("AE_SSE_HEARTBEAT_SECONDS", "12"))


async def _stream_with_heartbeat(stream_factory, heartbeat_s: float = None):
    """게이트웨이 이벤트 async-gen을 감싸 idle 시 heartbeat 이벤트를 주입한다.

    producer/consumer 큐 패턴 — producer 태스크가 원본 스트림을 소비해 큐에 넣고,
    consumer(본 제너레이터)는 heartbeat_s 내에 이벤트가 없으면 heartbeat를 yield한다.
    asyncio.wait_for로 producer를 취소하지 않으므로 원본 스트림이 끊기지 않는다.

    Args:
        stream_factory: 인자 없는 callable — 호출 시 async generator 반환.
        heartbeat_s: idle 임계(초). None이면 _SSE_HEARTBEAT_SECONDS.
    """
    hb = heartbeat_s if heartbeat_s is not None else _SSE_HEARTBEAT_SECONDS
    q: asyncio.Queue = asyncio.Queue()
    _SENTINEL = object()

    async def _pump():
        try:
            async for _e in stream_factory():
                await q.put(_e)
        except Exception as _ex:  # noqa: BLE001 — 스트림 예외를 이벤트로 전달
            await q.put({"type": "error", "message": f"stream error: {str(_ex)[:300]}"})
        finally:
            await q.put(_SENTINEL)

    task = asyncio.create_task(_pump())
    try:
        while True:
            try:
                evt = await asyncio.wait_for(q.get(), timeout=hb)
            except asyncio.TimeoutError:
                yield {"type": "heartbeat"}
                continue
            if evt is _SENTINEL:
                break
            yield evt
    finally:
        if not task.done():
            task.cancel()
            try:
                await task
            except BaseException:  # noqa: BLE001 — 취소/정리 예외 무시
                pass


def _resolve_callable_model_id(model_id, aws_profile, bedrock_user):
    """모델 ID를 실제 Bedrock 호출 가능한 형태로 변환.
    - ON_DEMAND only → prefix 제거 (prefix가 붙어있으면 떼어냄)
    - INFERENCE_PROFILE only → us. prefix 강제 (없으면 붙임)
    - 둘 다 / 알 수 없음 → prefix 있으면 유지, 없으면 us. 붙임 (기본값, 대부분 CRIS 커버)
    """
    if not model_id:
        return model_id
    # gateway-openai-models: OpenAI Responses 라우트 모델(openai.*)은 Bedrock
    # 추론 프로파일 대상이 아니다. us./eu. prefix를 붙이면 게이트웨이가 미지원으로
    # 거부하고 is_openai_model 판정도 깨진다. 원본 id를 그대로 반환한다(비침습).
    if is_openai_model(model_id):
        return model_id
    # Amazon 이미지 생성 모델(Nova Canvas / Titan Image)은 Bedrock ON_DEMAND
    # 이며 CRIS 추론 프로파일 대상이 아니다. list_foundation_models 캐시가
    # 비거나 이들을 누락하면 아래 기본 분기가 us. prefix를 붙여 Nova는
    # '/invoke allowlist 미스', Titan은 ResourceNotFoundException 을 일으킨다.
    # 원본(un-prefixed) id로 강제한다. (Stability 생성형은 인가 통과 상태라
    # 여기서 건드리지 않는다 — prefix 변경이 권한을 깨뜨릴 수 있음.)
    _img_raw = model_id
    for _ip in ("us.", "eu.", "global."):
        if _img_raw.startswith(_ip):
            _img_raw = _img_raw[len(_ip):]
            break
    if _img_raw in ("amazon.nova-canvas-v1:0", "amazon.titan-image-generator-v2:0"):
        return _img_raw
    # 이미 prefix 붙어있으면 원본 ID 추출
    raw_id = model_id
    for prefix in ("us.", "eu.", "global."):
        if model_id.startswith(prefix):
            raw_id = model_id[len(prefix):]
            break

    cache = _load_inference_types(aws_profile, bedrock_user)
    types = cache.get(raw_id, [])

    has_on_demand = "ON_DEMAND" in types
    has_inference_profile = "INFERENCE_PROFILE" in types

    if has_on_demand and not has_inference_profile:
        # ON_DEMAND 전용 — prefix 제거
        return raw_id
    if has_inference_profile and not has_on_demand:
        # INFERENCE_PROFILE 전용 — us. prefix 강제
        if not any(model_id.startswith(p) for p in ("us.", "eu.", "global.")):
            return f"us.{raw_id}"
        return model_id
    # 둘 다 또는 알 수 없음 → 기본 CRIS (기존 동작 유지)
    if not any(model_id.startswith(p) for p in ("us.", "eu.", "global.")):
        return f"us.{raw_id}"
    return model_id


def _is_expired_error(result):
    """응답이 토큰 만료 에러인지 판단."""
    err = ""
    if isinstance(result, dict):
        err = result.get("error", "")
    elif isinstance(result, str):
        err = result
    low = err.lower()
    return "expired" in low or "security token" in low


async def _refresh_and_retry_gw(gw, aws_profile, bedrock_user):
    """토큰 만료 시 자격증명을 다시 assume role하여 주입."""
    try:
        import boto3 as b3
        from botocore.credentials import Credentials as BotoCreds
        # boto3 세션 캐시 초기화
        b3.DEFAULT_SESSION = None
        gw.force_refresh_creds()
        session = b3.Session(profile_name=aws_profile)
        sts = session.client("sts")
        account = sts.get_caller_identity()["Account"]
        if bedrock_user:
            assumed = sts.assume_role(
                RoleArn=f"arn:aws:iam::{account}:role/BedrockUser-{bedrock_user}",
                RoleSessionName="ai-editor-refresh",
            )
            c = assumed["Credentials"]
            gw.inject_credentials(c["AccessKeyId"], c["SecretAccessKey"], c["SessionToken"])
            print(f"[AutoRefresh] BedrockUser-{bedrock_user} 자격증명 재주입 성공")
            return True
        else:
            fc = session.get_credentials().get_frozen_credentials()
            gw.inject_credentials(fc.access_key, fc.secret_key, fc.token)
            print(f"[AutoRefresh] 프로파일 {aws_profile} 자격증명 재주입 성공")
            return True
    except Exception as e:
        print(f"[AutoRefresh] 자격증명 재주입 실패: {e}")
        return False


@app.api_route("/health", methods=["GET", "HEAD"])
async def health():
    return {
        "status": "ok",
        "service": "ai-editor-engine",
        "timestamp": datetime.utcnow().isoformat(),
        "version": __version__,
    }



@app.get("/api/debug/cwd")
async def debug_cwd():
    """Return server cwd + the actual root used for .generated/ files.

    `generatedRoot`는 _resolve_local_root이 사용하는 실제 쓰기 루트.
    file-preview-panel이 이 값을 보고 .generated/ 위치를 정확히 찾는다.
    """
    return {
        "cwd": os.getcwd(),
        "generatedRoot": _resolve_local_root(""),
        "envRoot": os.environ.get("AE_GENERATED_ROOT", ""),
    }


@app.get("/api/debug/image-gen-status")
async def debug_image_gen_status():
    """Diagnostic surface for image-generation routing.

    Returns the circuit-breaker state, the configured `IMAGE_MODELS` chain,
    a `_select_image_models` preview for a representative prompt, the
    fresh-read values of the four image-gen environment variables, and the
    last 10 attempt outcomes from `_IMAGE_GEN_ATTEMPTS`.

    Always returns HTTP 200 — the diagnostic must NOT be gated by the
    breaker (operator needs to see the state precisely WHEN the breaker
    is closed).

    Spec: media-output-quality (bugfix) — Property 4 / Req 1.4, 2.4.
    """
    import time as _t
    now = _t.time()
    disabled_at = _IMAGE_GEN_CIRCUIT.get("disabled_at", 0) or 0
    ttl = _IMAGE_GEN_CIRCUIT.get("ttl", 300) or 300
    ttl_remaining = max(0.0, ttl - (now - disabled_at)) if disabled_at > 0 else 0.0
    is_broken = _image_gen_is_circuit_broken()
    # _select_image_models takes (prompt, hint=""), not (prompt, None) — pass
    # an empty string so the call signature is honored.
    try:
        select_preview = _select_image_models("test architecture diagram", "")
    except Exception as _e:  # never let the diagnostic crash on a stub
        select_preview = []
    # Vertex 이미지(최고 품질 폴백) 활성 상태 — 로그인 시 자동 해석된 싱글톤 반영.
    _vertex_status = {"enabled": False, "projectId": None, "reason": "not-initialized"}
    try:
        from ai_engine import vertex_image_module as _vmod
        _vc = _vmod.get_vertex_image_client()
        _vertex_status["enabled"] = bool(getattr(_vc, "enabled", False))
        _vertex_status["projectId"] = getattr(_vc, "_project_id", None)
        _vertex_status["googleAuthInstalled"] = bool(getattr(_vmod, "_GOOGLE_AUTH_AVAILABLE", False))
        if _vertex_status["enabled"]:
            _vertex_status["reason"] = "ok"
        elif not _vertex_status["googleAuthInstalled"]:
            _vertex_status["reason"] = "google-auth/requests 미설치 — pip install google-auth requests"
        else:
            _vertex_status["reason"] = "키 미해석(Secrets Manager 권한/키 확인) — 로그인 후 재시도"
    except Exception as _ve:
        _vertex_status["reason"] = str(_ve)[:160]
    return {
        "circuit": {
            "disabled_at": float(disabled_at),
            "ttl": int(ttl),
            "ttlRemainingSec": float(ttl_remaining),
            "isBroken": bool(is_broken),
        },
        "models": list(IMAGE_MODELS),
        "selectPreview": list(select_preview),
        "vertex": _vertex_status,
        "env": {
            "AE_IMAGE_PARALLEL_N": os.getenv("AE_IMAGE_PARALLEL_N"),
            "AE_IMAGE_QUALITY_THRESHOLD": os.getenv("AE_IMAGE_QUALITY_THRESHOLD"),
            "AE_FORCE_NATIVE_DIAGRAM": os.getenv("AE_FORCE_NATIVE_DIAGRAM"),
            "AE_DISABLE_HTML_SLIDES": os.getenv("AE_DISABLE_HTML_SLIDES"),
            "AE_ENABLE_HTML_SLIDES": os.getenv("AE_ENABLE_HTML_SLIDES"),
            "AE_DISABLE_VERTEX_IMAGE": os.getenv("AE_DISABLE_VERTEX_IMAGE"),
        },
        "recentAttempts": list(_IMAGE_GEN_ATTEMPTS),
    }


@app.get("/api/debug/bridge")
async def debug_bridge():
    """Debug: show bridge state."""
    _refresh_bridge_discovery()
    return {
        "bridge_url": _BRIDGE_URL,
        "bridge_token_set": bool(_BRIDGE_TOKEN),
        "is_remote": _bridge_is_remote() if _BRIDGE_URL else False,
    }


@app.get("/api/debug/openai-test")
async def debug_openai_test(request: Request):
    """gateway-openai-models 진단 — OpenAI 라우트 원응답을 그대로 노출.

    우리 쪽 재시도/예외 래핑을 모두 우회하고 저수준 호출(_openai_request_blocking)을
    한 번 수행해 게이트웨이가 반환한 상태코드/본문을 그대로 반환한다.
    403/422/500 + 본문으로 실제 원인을 즉시 식별할 수 있다.

    사용: GET /api/debug/openai-test?model=openai.gpt-5.5&profile=<프로파일>
    """
    qp = request.query_params
    model = qp.get("model", "openai.gpt-5.5")
    profile = qp.get("profile", os.environ.get("AWS_PROFILE", "default"))
    bedrock_user = qp.get("bedrockUser", os.environ.get("AE_BEDROCK_USER", ""))
    prompt = qp.get("prompt", "Say 'pong' in one word.")

    from ai_engine.gateway_module import mask_token

    gw = _get_gw(profile, bedrock_user)
    messages = [{"role": "user", "content": prompt}]
    result: dict = {
        "model_requested": model,
        "is_openai_model": is_openai_model(model),
        "gateway_url": gw.gateway_url,
        "profile": profile,
        "bedrock_user": bedrock_user or "(none)",
    }
    try:
        body = gw._build_openai_payload(model, messages, system_prompt="")
        body_bytes = json.dumps(body).encode()
        result["sent_body"] = body
        # 서명 헤더 — 토큰류는 마스킹해서 어떤 헤더가 붙는지만 노출
        try:
            url = f"{gw.gateway_url}/openai/responses"
            headers = gw._sign("POST", url, body_bytes)
            result["signed_header_keys"] = sorted(headers.keys())
            if "Authorization" in headers:
                result["authorization_preview"] = mask_token(headers["Authorization"])
        except Exception as _se:
            result["sign_error"] = str(_se)[:300]
            return JSONResponse(result, status_code=200)

        loop = asyncio.get_event_loop()
        raw = await loop.run_in_executor(
            None, gw._openai_request_blocking, "POST", url, body_bytes, 30
        )
        # 원응답 그대로(본문은 1500자까지)
        result["sync_status"] = raw.get("status")
        body_text = raw.get("body", "") or ""
        result["sync_body"] = body_text[:1500]
        if raw.get("error"):
            result["sync_error"] = str(raw.get("error"))[:500]

        # 실제 채팅 경로(route_openai_chat → to_converse) 전체 실행 — 어댑터 검증
        try:
            _conv = await route_openai_chat(gw, model, messages, system_prompt="", timeout=60)
            _txt = ""
            for _c in _conv.get("output", {}).get("message", {}).get("content", []):
                if isinstance(_c, dict) and _c.get("text"):
                    _txt += _c["text"]
            result["route_ok"] = True
            result["route_text"] = _txt[:500]
            result["route_usage"] = _conv.get("usage")
        except Exception as _re:
            import traceback as _tb
            result["route_ok"] = False
            result["route_error"] = f"{type(_re).__name__}: {str(_re)[:500]}"
            result["route_traceback"] = _tb.format_exc()[-1200:]

        # job 라우트도 한 번 시험(제출만)
        try:
            jurl = f"{gw.gateway_url}/openai/responses-jobs"
            jheaders = gw._sign("POST", jurl, body_bytes)
            jraw = await loop.run_in_executor(
                None, gw._openai_request_blocking, "POST", jurl, body_bytes, 30
            )
            result["jobs_status"] = jraw.get("status")
            result["jobs_body"] = (jraw.get("body", "") or "")[:800]
        except Exception as _je:
            result["jobs_error"] = str(_je)[:300]

        # 도구(함수호출) 지원 확인 — tools=1 일 때만. GPT가 function_call을
        # 반환하는지/그 형식을 그대로 노출해 tool 루프 구현 근거를 확보한다.
        if qp.get("tools") in ("1", "true", "yes"):
            try:
                tools = _agent_tools_to_openai()
                # 도구를 유발하는 입력으로 교체(기본 'pong' 프롬프트는 도구를 안 부름)
                tool_prompt = qp.get("toolPrompt", "현재 폴더에서 README 파일을 찾아줘. search_files 도구를 사용해.")
                tinput = [{"role": "user", "content": [{"type": "input_text", "text": tool_prompt}]}]
                tbody = {"model": model, "input": tinput, "tools": tools, "tool_choice": "auto"}
                tbody_bytes = json.dumps(tbody).encode()
                turl = f"{gw.gateway_url}/openai/responses"
                traw = await loop.run_in_executor(
                    None, gw._openai_request_blocking, "POST", turl, tbody_bytes, 60
                )
                result["tools_status"] = traw.get("status")
                result["tools_body"] = (traw.get("body", "") or "")[:2500]
            except Exception as _te:
                result["tools_error"] = str(_te)[:400]
    except Exception as e:
        result["exception"] = f"{type(e).__name__}: {str(e)[:500]}"
    return JSONResponse(result, status_code=200)


@app.post("/api/attachments/extract-zip")
async def extract_zip_attachment(request: Request):
    """Receive a zip file via multipart upload, extract to a sandboxed temp dir,
    and return a structured listing of internal files for chat context.

    Security:
      - Refuses absolute paths and any member containing '..' segments.
      - Per-file 5MB cap, 50MB total cap, max 200 files.
      - Each upload gets its own tempfile.mkdtemp() dir (isolation).

    Response:
      {
        "extractDir": str (absolute),
        "files": [{"path": str (rel), "size": int, "type": "text"|"binary"|"image"}],
        "totalFiles": int,
        "totalBytes": int,
        "skipped": [str, ...]
      }
    """
    import zipfile
    import tempfile

    try:
        form = await request.form()
    except Exception as e:
        return JSONResponse({"error": f"bad-form: {str(e)[:200]}"}, status_code=400)

    upload = form.get("file")
    if upload is None or not hasattr(upload, "read"):
        return JSONResponse({"error": "no-file"}, status_code=400)

    extract_root = tempfile.mkdtemp(prefix="ae_zip_extract_")
    zip_path = os.path.join(extract_root, "_input.zip")
    try:
        data = await upload.read()
        with open(zip_path, "wb") as f:
            f.write(data)
    except Exception as e:
        return JSONResponse({"error": f"write-failed: {str(e)[:200]}"}, status_code=500)

    files_listing = []
    skipped = []
    total_bytes = 0
    MAX_FILES = 200
    MAX_FILE_SIZE = 5 * 1024 * 1024  # 5 MB per file
    MAX_TOTAL_BYTES = 50 * 1024 * 1024  # 50 MB total

    TEXT_EXTS = {
        ".py", ".js", ".ts", ".tsx", ".jsx", ".html", ".css", ".scss",
        ".json", ".yaml", ".yml", ".md", ".txt", ".sh", ".bash",
        ".sql", ".go", ".rs", ".java", ".cpp", ".c", ".h", ".hpp",
        ".xml", ".csv", ".toml", ".ini", ".cfg", ".env",
    }
    IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}

    try:
        with zipfile.ZipFile(zip_path, "r") as zf:
            for member in zf.namelist():
                # Path-traversal protection: refuse absolute paths and '..'
                if member.startswith("/") or member.startswith("\\"):
                    skipped.append(f"{member} (absolute path)")
                    continue
                parts = member.replace("\\", "/").split("/")
                if any(p == ".." for p in parts):
                    skipped.append(f"{member} (parent traversal)")
                    continue

                info = zf.getinfo(member)
                if info.is_dir():
                    continue

                if len(files_listing) >= MAX_FILES:
                    skipped.append(f"{member} (max files reached)")
                    continue
                if info.file_size > MAX_FILE_SIZE:
                    skipped.append(f"{member} (>5MB)")
                    continue
                if total_bytes + info.file_size > MAX_TOTAL_BYTES:
                    skipped.append(f"{member} (total cap)")
                    continue

                ext = os.path.splitext(member)[1].lower()
                if ext in TEXT_EXTS:
                    file_type = "text"
                elif ext in IMAGE_EXTS:
                    file_type = "image"
                else:
                    file_type = "binary"

                # Extract this single member into the sandbox
                zf.extract(member, extract_root)
                files_listing.append({
                    "path": member,
                    "size": int(info.file_size),
                    "type": file_type,
                })
                total_bytes += info.file_size

        # Drop the original zip from the sandbox so the listing is clean
        try:
            os.remove(zip_path)
        except Exception:
            pass

        return {
            "extractDir": extract_root,
            "files": files_listing,
            "totalFiles": len(files_listing),
            "totalBytes": total_bytes,
            "skipped": skipped,
        }
    except zipfile.BadZipFile:
        return JSONResponse({"error": "bad-zip"}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)[:200]}, status_code=500)


@app.post("/api/reset-cache")
async def reset_cache(request: Request):
    """Gateway 클라이언트 캐시 초기화 + 선택적 자격증명 주입."""
    _gw_cache.clear()
    _quota_cache["used_krw"] = 0
    _quota_cache["remaining_krw"] = 0
    _quota_cache["limit_krw"] = 0
    _quota_cache["last_updated"] = ""
    _quota_cache["user"] = ""
    try:
        import boto3
        boto3.DEFAULT_SESSION = None
    except Exception:
        pass
    # 자격증명 직접 주입 (Electron에서 전달)
    try:
        body = await request.json()
        creds = body.get("credentials")
        if creds and creds.get("AWS_ACCESS_KEY_ID"):
            profile = body.get("profile", "bedrock-gw")
            user = body.get("bedrockUser", "")
            # 새 GatewayClient 생성 후 자격증명 주입
            from ai_engine.gateway_module import GatewayClient
            gw = GatewayClient(
                gateway_url=os.environ.get("GATEWAY_URL", "https://5l764dh7y9.execute-api.us-west-2.amazonaws.com/v1"),
                aws_profile=profile,
                region=os.environ.get("AWS_REGION", "us-west-2"),
                bedrock_user=user,
            )
            # SSO 기본 자격증명으로 BedrockUser assume role 시도
            try:
                import boto3 as b3
                from botocore.credentials import Credentials as BotoCreds
                # 전달받은 SSO 자격증명으로 임시 세션 생성
                tmp_session = b3.Session()
                sts = tmp_session.client(
                    "sts",
                    aws_access_key_id=creds["AWS_ACCESS_KEY_ID"],
                    aws_secret_access_key=creds["AWS_SECRET_ACCESS_KEY"],
                    aws_session_token=creds.get("AWS_SESSION_TOKEN", ""),
                    region_name=creds.get("AWS_DEFAULT_REGION", "us-west-2"),
                )
                account = sts.get_caller_identity()["Account"]
                if user:
                    assumed = sts.assume_role(
                        RoleArn=f"arn:aws:iam::{account}:role/BedrockUser-{user}",
                        RoleSessionName="ai-editor",
                    )
                    c = assumed["Credentials"]
                    gw.inject_credentials(c["AccessKeyId"], c["SecretAccessKey"], c["SessionToken"])
                    print(f"[Cache] BedrockUser-{user} assume role 성공")
                else:
                    gw.inject_credentials(
                        creds["AWS_ACCESS_KEY_ID"],
                        creds["AWS_SECRET_ACCESS_KEY"],
                        creds.get("AWS_SESSION_TOKEN", ""),
                    )
                key = f"{profile}:{user}"
                _gw_cache[key] = gw
            except Exception as e:
                print(f"[Cache] assume role 실패: {e}")

            # === Vertex 이미지 자동 활성화 (모든 사용자 무설정 적용) ===
            # 로그인 시 주입된 SSO 자격증명으로 AWS Secrets Manager에서 GCP 키를
            # 자동 해석한다(secretsmanager:GetSecretValue 권한 보유 사용자). 성공 시
            # 로컬 캐시에 저장돼 이후 오프라인에서도 동작. 실패해도 로그인은 정상 진행.
            try:
                from ai_engine.vertex_image_module import (
                    reset_vertex_image_client, get_vertex_image_client,
                )
                reset_vertex_image_client()
                _vx_creds = {
                    "accessKeyId": creds["AWS_ACCESS_KEY_ID"],
                    "secretAccessKey": creds["AWS_SECRET_ACCESS_KEY"],
                    "sessionToken": creds.get("AWS_SESSION_TOKEN", ""),
                    "region": creds.get("AWS_DEFAULT_REGION", "us-west-2"),
                }
                _vx = get_vertex_image_client(aws_profile=profile, credentials=_vx_creds)
                if getattr(_vx, "enabled", False):
                    print("[Cache] Vertex 이미지 자동 활성화 성공 (Secrets Manager 키 해석)")
                else:
                    print("[Cache] Vertex 비활성 — Secrets Manager 키 없음/권한 없음(흐름도는 네이티브로 동작)")
            except Exception as _ve:
                print(f"[Cache] Vertex 자동 활성화 시도 실패(무시): {str(_ve)[:160]}")
    except Exception:
        pass
    return {"status": "ok", "message": "cache cleared"}


@app.post("/api/rag/index")
async def rag_index(request: Request):
    """프로젝트 인덱싱 수동 트리거."""
    body = await request.json()
    project_path = body.get("projectPath", "")
    if not project_path or not os.path.isdir(project_path):
        return JSONResponse(content={"error": "Invalid project path"}, status_code=400)
    from ai_engine.rag.context_builder import get_indexer
    idx = get_indexer(project_path)
    count = idx.index_project(project_path)
    return {"status": "ok", "chunks": count, "files": len(set(c.file_path for c in idx.chunks))}


@app.get("/api/rag/status")
async def rag_status(request: Request):
    """RAG 인덱싱 상태 조회."""
    project_path = request.query_params.get("projectPath", "")
    if not project_path:
        return {"indexed": False, "chunks": 0}
    from ai_engine.rag.context_builder import _indexer_cache
    if project_path in _indexer_cache:
        idx = _indexer_cache[project_path]
        return {"indexed": True, "chunks": len(idx.chunks), "files": len(set(c.file_path for c in idx.chunks))}
    return {"indexed": False, "chunks": 0}


@app.get("/api/models")
@app.post("/api/models")
async def list_models(request: Request):
    """Return available models. POST로 자격증명을 직접 전달 가능."""
    profile = request.query_params.get("profile", os.environ.get("AWS_PROFILE", "default"))
    
    # POST body에서 자격증명 직접 받기
    creds_override = None
    if request.method == "POST":
        try:
            body = await request.json()
            if body.get("accessKeyId"):
                creds_override = body
                profile = body.get("profile", profile)
        except Exception:
            pass
    
    try:
        import boto3

        if creds_override:
            # 전달받은 자격증명으로 직접 클라이언트 생성
            client = boto3.client(
                "bedrock",
                aws_access_key_id=creds_override["accessKeyId"],
                aws_secret_access_key=creds_override["secretAccessKey"],
                aws_session_token=creds_override.get("sessionToken", ""),
                region_name=creds_override.get("region", os.environ.get("AWS_REGION", "us-west-2")),
            )
        else:
            session = boto3.Session(
                profile_name=profile,
                region_name=os.environ.get("AWS_REGION", "us-west-2"),
            )
            client = session.client("bedrock")
        
        # list_foundation_models는 등록된 전체 카탈로그를 반환하지만 모든 모델이
        # converse API로 호출 가능한 건 아님. 실제 호출 가능한 모델만 걸러내려면
        # list_inference_profiles (CRIS) 결과와 교차 검증 필요.
        # - Cross-region Inference Profile이 있는 모델만 converse 호출 가능
        # - 없는 모델은 ValidationException 발생
        try:
            profiles_resp = client.list_inference_profiles()
            callable_profile_ids = set()
            profile_id_to_name = {}
            for p in profiles_resp.get("inferenceProfileSummaries", []):
                if p.get("status", "").upper() != "ACTIVE":
                    continue
                pid = p.get("inferenceProfileId", "")
                pname = p.get("inferenceProfileName", pid)
                if pid:
                    callable_profile_ids.add(pid)
                    profile_id_to_name[pid] = pname
        except Exception as _e:
            # 프로파일 조회 실패 → 폴백으로 foundation model 전체 사용 (기존 동작)
            print(f"[Models] list_inference_profiles 실패, 폴백: {_e}")
            callable_profile_ids = None
            profile_id_to_name = {}

        resp = client.list_foundation_models()
        catalog = {}
        seen_model_keys = set()  # 리전/컨텍스트 변형 중복 제거

        skip_output = []  # IMAGE/VIDEO/EMBEDDING 별도 카탈로그로 분리
        image_catalog = {}
        seen_image_keys = set()
        video_catalog = {}
        seen_video_keys = set()
        embed_catalog = {}
        seen_embed_keys = set()
        rerank_catalog = {}
        seen_rerank_keys = set()
        for m in resp.get("modelSummaries", []):
            modes = m.get("outputModalities", [])
            if m.get("modelLifecycle", {}).get("status") in ["EOL"]:
                continue
            input_modes = m.get("inputModalities", [])

            _mid_raw = m["modelId"]
            _no_region = _mid_raw
            for _pfx in ("us.", "eu.", "global."):
                if _no_region.startswith(_pfx):
                    _no_region = _no_region[len(_pfx):]
                    break

            # === Embedding models ===
            if "EMBEDDING" in modes or "embed" in _no_region.lower():
                if _no_region in seen_embed_keys:
                    continue
                seen_embed_keys.add(_no_region)
                inference_types = m.get("inferenceTypesSupported", [])
                if inference_types and "ON_DEMAND" not in inference_types and "INFERENCE_PROFILE" not in inference_types:
                    continue
                callable_id = _mid_raw
                if callable_profile_ids is not None and "ON_DEMAND" not in inference_types:
                    for pid in callable_profile_ids:
                        if pid.endswith(_mid_raw) or pid.endswith(f"{_mid_raw}:0") or _mid_raw in pid:
                            callable_id = pid
                            break
                _prov = m.get("providerName", "Unknown")
                if _prov not in embed_catalog:
                    embed_catalog[_prov] = []
                embed_catalog[_prov].append({"id": callable_id, "name": m.get("modelName", _mid_raw)})
                continue

            # === Rerank models ===
            if "rerank" in _no_region.lower():
                if _no_region in seen_rerank_keys:
                    continue
                seen_rerank_keys.add(_no_region)
                inference_types = m.get("inferenceTypesSupported", [])
                if inference_types and "ON_DEMAND" not in inference_types and "INFERENCE_PROFILE" not in inference_types:
                    continue
                callable_id = _mid_raw
                if callable_profile_ids is not None and "ON_DEMAND" not in inference_types:
                    for pid in callable_profile_ids:
                        if pid.endswith(_mid_raw) or pid.endswith(f"{_mid_raw}:0") or _mid_raw in pid:
                            callable_id = pid
                            break
                _prov = m.get("providerName", "Unknown")
                if _prov not in rerank_catalog:
                    rerank_catalog[_prov] = []
                rerank_catalog[_prov].append({"id": callable_id, "name": m.get("modelName", _mid_raw)})
                continue

            # === Video generation models ===
            if "VIDEO" in modes:
                if _no_region in seen_video_keys:
                    continue
                seen_video_keys.add(_no_region)
                inference_types = m.get("inferenceTypesSupported", [])
                if inference_types and "ON_DEMAND" not in inference_types and "INFERENCE_PROFILE" not in inference_types:
                    continue
                callable_id = _mid_raw
                if callable_profile_ids is not None and "ON_DEMAND" not in inference_types:
                    for pid in callable_profile_ids:
                        if pid.endswith(_mid_raw) or pid.endswith(f"{_mid_raw}:0") or _mid_raw in pid:
                            callable_id = pid
                            break
                _prov = m.get("providerName", "Unknown")
                if _prov not in video_catalog:
                    video_catalog[_prov] = []
                video_catalog[_prov].append({"id": callable_id, "name": m.get("modelName", _mid_raw)})
                continue

            # === Image generation/edit models ===
            if "IMAGE" in modes:
                _mid = m["modelId"]
                _parts = _no_region.split(":")
                _img_base = _parts[0] + ":" + _parts[1] if len(_parts) >= 2 else _no_region
                if _img_base in seen_image_keys:
                    continue
                seen_image_keys.add(_img_base)
                inference_types = m.get("inferenceTypesSupported", [])
                callable_id = _mid
                if callable_profile_ids is not None and "ON_DEMAND" not in inference_types:
                    for pid in callable_profile_ids:
                        if pid.endswith(_mid) or pid.endswith(f"{_mid}:0") or _mid in pid:
                            callable_id = pid
                            break
                _prov = m.get("providerName", "Unknown")
                if _prov not in image_catalog:
                    image_catalog[_prov] = []
                image_catalog[_prov].append({
                    "id": callable_id,
                    "name": m.get("modelName", _mid),
                })
                continue

            # === Text models ===
            if "TEXT" not in input_modes:
                continue
            if "TEXT" not in input_modes:
                continue
            inference_types = m.get("inferenceTypesSupported", [])
            if inference_types and "ON_DEMAND" not in inference_types and "INFERENCE_PROFILE" not in inference_types:
                continue
            if m.get("responseStreamingSupported") is False:
                continue

            # 실제 호출 가능 여부: CRIS profile 존재 여부로 판단
            model_id = m["modelId"]
            # 중복 제거: 리전 prefix(us./eu./global.) + context window 변형(:8k,:20k,:1000k,:mm) 정규화
            _no_region = model_id
            for _pfx in ("us.", "eu.", "global."):
                if _no_region.startswith(_pfx):
                    _no_region = _no_region[len(_pfx):]
                    break
            _parts = _no_region.split(":")
            if len(_parts) >= 2:
                _base_key = _parts[0] + ":" + _parts[1]
            else:
                _base_key = _no_region
            if _base_key in seen_model_keys:
                continue
            seen_model_keys.add(_base_key)
            if len(_parts) > 2:
                model_id = _parts[0] + ":" + _parts[1]  # 기본 ID(:0) 우선

            if callable_profile_ids is not None and "INFERENCE_PROFILE" in inference_types and "ON_DEMAND" not in inference_types:
                # 이 모델은 CRIS profile 필수 — profile 존재 여부 확인
                has_profile = any(
                    pid.endswith(model_id) or pid.endswith(f"{model_id}:0") or model_id in pid
                    for pid in callable_profile_ids
                )
                if not has_profile:
                    continue  # CRIS profile 없으면 호출 불가 → 스킵

            provider = m.get("providerName", "Unknown")
            if provider not in catalog:
                catalog[provider] = []
            catalog[provider].append({
                "id": model_id,
                "name": m.get("modelName", model_id),
            })
        # 라우팅 매트릭스가 동적으로 카탈로그를 확장 후보로 사용하도록 캐시 업데이트.
        # 이렇게 해야 게이트웨이에 Gemma 등 새 모델이 추가되면 코드 변경 없이도 라우팅
        # 후보로 자동 활용된다.
        try:
            _update_gateway_model_cache(catalog)
        except Exception as _ce:
            print(f"[Models] cache update 실패: {_ce}")
        # === OpenAI 모델 병합 (gateway-openai-models, 요구사항 1) ===
        # 게이트웨이의 OpenAI Responses 라우트 모델을 Bedrock 카탈로그에 병합한다.
        # try/except graceful — 실패 시 원인 로그 후 Bedrock-only 반환(요구사항 1.7).
        # OpenAI 0개면 catalog가 변경되지 않아 baseline 바이트 보존(요구사항 8.1).
        try:
            from ai_engine import openai_catalog as _oc
            _src = _oc.get_catalog_source({})
            _openai_entries = _src.list_models()
            catalog = _oc.merge_openai_into_catalog(catalog, _openai_entries)
        except Exception as _oe:
            print(f"[Models] OpenAI 병합 실패 → Bedrock-only: {str(_oe)[:200]}")
        _text_count = sum(len(v) for v in catalog.values())
        _image_count = sum(len(v) for v in image_catalog.values())
        _video_count = sum(len(v) for v in video_catalog.values())
        _embed_count = sum(len(v) for v in embed_catalog.values())
        _rerank_count = sum(len(v) for v in rerank_catalog.values())
        return JSONResponse(content={
            "models": catalog,
            "image_models": image_catalog,
            "video_models": video_catalog,
            "embed_models": embed_catalog,
            "rerank_models": rerank_catalog,
            # `count`는 프런트 드롭다운(ALL_MODELS)이 실제로 병합해 보여주는
            # 전체 호출가능 모델 수(text+image+video+embed+rerank)와 일치시킨다.
            # 과거엔 text-only(예: 74)만 반환해 실제 카탈로그(~100)와 불일치했다.
            "count": _text_count + _image_count + _video_count + _embed_count + _rerank_count,
            "text_count": _text_count,
            "image_count": _image_count,
            "video_count": _video_count,
            "embed_count": _embed_count,
            "rerank_count": _rerank_count,
        })
    except Exception as e:
        return JSONResponse(content={"models": {}, "error": str(e)})


# ─────────────────────────────────────────────────────────────────
# Intent Classifier — LLM 기반 의도 분류
# 사용자 메시지를 분석하여 최적 실행 전략을 결정한다.
# ─────────────────────────────────────────────────────────────────

INTENT_CLASSIFIER_PROMPT = """사용자 요청을 분석하여 최적 실행 전략을 결정하세요.

[분류 카테고리]
- file_generation: 실제 파일을 생성/수정해야 함 (PDF/XLSX/DOCX/PPTX/HWP/이미지/코드파일)
- code_change: 기존 코드 파일을 수정/리팩토링/디버깅
- analysis: 코드/문서/데이터 분석, 설명, 리뷰
- generation_text: 글/콘텐츠/아이디어 생성 (파일 저장 불필요)
- reasoning: 수학/논리/추론 문제 (단계별 사고 필요, 도구 X)
- translation: 번역, 다국어 변환
- long_context: 매우 긴 문서/로그/대화 처리 (>10K 자)
- simple_qa: 간단한 질문, 사실 확인, 짧은 응답

[복잡도]
- simple: 1-2 문단 응답으로 충분
- moderate: 여러 단계 추론 필요
- complex: 다중 파일/모듈/긴 작업

[병렬 유용성]
- true: 여러 모델의 다른 관점 비교가 가치 있음 (분석, 창작, 리뷰)
- false: 정답이 하나거나 도구 실행이 필수 (파일 생성, 코드 수정)

반드시 아래 JSON 형식으로만 응답하세요 (마크다운 없이):
{
  "intent": "file_generation" | "code_change" | "analysis" | "generation_text" | "reasoning" | "translation" | "long_context" | "simple_qa",
  "needs_tools": true | false,
  "complexity": "simple" | "moderate" | "complex",
  "parallel_useful": true | false,
  "file_types": ["pdf", "xlsx", ...],
  "reasoning": "한 문장 이유"
}
"""


# intent → task_type 매핑 — _specialized_model_for_task 인자로 전달.
# 이 매핑이 다양한 모델 활용의 핵심: 각 intent가 적합한 모델 풀로 라우팅됨.
_INTENT_TO_TASK = {
    "file_generation": "file_generation",   # → Sonnet (도구 호출)
    "code_change":     "code",              # → Opus
    "analysis":        "doc_analysis",      # → Qwen3 / Sonnet
    "generation_text": "general_chat",      # → Nova Lite / Llama / Haiku
    "reasoning":       "reasoning",         # → DeepSeek-R1 / Qwen3
    "translation":     "translation",       # → Qwen3 (한국어 강함)
    "long_context":    "long_context",      # → Qwen3-235B
    "simple_qa":       "simple_qa",         # → Nova Lite / Haiku
}


@app.post("/api/agents/classify-intent")
async def classify_intent(request: Request):
    """사용자 메시지의 의도를 LLM으로 분류한다. Haiku로 빠르게(<1초)."""
    body = await request.json()
    prompt = body.get("prompt", "")
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))

    if not prompt or len(prompt.strip()) < 2:
        return JSONResponse({
            "intent": "simple_qa", "needs_tools": False,
            "complexity": "simple", "parallel_useful": False,
            "file_types": [], "reasoning": "empty or trivial prompt"
        })

    gw = _get_gw(aws_profile, bedrock_user)
    # task별 자동 라우팅 — Haiku 우선(빠름/저렴)이되, 특정 Haiku가
    # 게이트웨이에서 응답이 느리거나 불가할 수 있으므로(예: haiku-4-5
    # 무응답 → 10s 하드컷 → 항상 degraded 되던 문제), 확정 활성 모델
    # (Sonnet, file_generation tier)로 폴백하는 후보 체인을 순차 시도한다.
    _primary_cls = _specialized_model_for_task(
        "intent_classifier", None,
        aws_profile=aws_profile, bedrock_user=bedrock_user,
    )
    _fallback_cls = _specialized_model_for_task(
        "file_generation", None,
        aws_profile=aws_profile, bedrock_user=bedrock_user,
    )
    _cls_candidates = []
    for _m in (_primary_cls, _fallback_cls):
        if _m and _m not in _cls_candidates:
            _cls_candidates.append(_m)
    # 후보별 타임아웃 — 비스트리밍 converse는 게이트웨이에서 ~40-50s로
    # 매우 느려 분류가 항상 타임아웃됐다. 빠른 스트리밍 경로(~8s 첫 토큰)로
    # 호출하고 현실적 타임아웃을 둔다.
    _cls_timeouts = [25, 35]

    try:
        result = None
        _cls_last_err = None
        for _idx, _cm in enumerate(_cls_candidates):
            _to = _cls_timeouts[_idx] if _idx < len(_cls_timeouts) else 35
            try:
                _r = await asyncio.wait_for(
                    gw.converse_stream_live(
                        model_id=_cm,
                        messages=[{"role": "user", "content": [{"text": prompt[:1500]}]}],
                        system_prompt=INTENT_CLASSIFIER_PROMPT,
                    ),
                    timeout=_to,
                )
            except Exception as _ce:
                _cls_last_err = _ce
                print(f"[Intent] 후보 {_cm} 실패({type(_ce).__name__ or 'timeout'}) → 다음 후보 시도")
                continue
            if _r.get("decision") == "ALLOW":
                result = _r
                break
            _cls_last_err = RuntimeError(f"classifier failed: {_r.get('error') or _r.get('decision')}")
        if result is None:
            raise _cls_last_err or RuntimeError("classifier: 모든 후보 실패")

        output = result.get("output", {}).get("message", {}).get("content", [])
        text = "\n".join(c.get("text", "") for c in output if "text" in c).strip()
        # JSON 추출
        m = re.search(r"\{[\s\S]*\}", text)
        if not m:
            raise RuntimeError(f"no JSON in classifier response: {text[:200]}")
        parsed = json.loads(m.group(0))
        # 기본값 보정
        return JSONResponse({
            "intent": parsed.get("intent", "simple_qa"),
            "needs_tools": bool(parsed.get("needs_tools", False)),
            "complexity": parsed.get("complexity", "simple"),
            "parallel_useful": bool(parsed.get("parallel_useful", False)),
            "file_types": parsed.get("file_types", []) or [],
            "reasoning": parsed.get("reasoning", "")[:200],
        })
    except Exception as e:
        # 분류 실패 시 안전한 기본값 (단순 QA로 처리). 어떤 경우에도 200 + 유효 JSON 반환.
        # degraded=True로 표시하고, SSO/자격증명 만료를 감지하면 needsReauth를 세워
        # 프론트가 사용자에게 재로그인을 안내할 수 있게 한다(분류 자체는 fallback으로 진행).
        _emsg = str(e)
        print(f"[Intent] 분류 실패: {_emsg[:200]}")
        _needs_reauth = _is_expired_error(_emsg) if "_is_expired_error" in globals() else (
            any(t in _emsg.lower() for t in (
                "sso session", "expired", "token", "credential", "unable to locate credentials",
                "the security token", "expiredtoken",
            ))
        )
        return JSONResponse({
            "intent": "simple_qa", "needs_tools": False,
            "complexity": "simple", "parallel_useful": False,
            "file_types": [], "reasoning": f"classifier failed: {_emsg[:100]}",
            "degraded": True,
            "needsReauth": bool(_needs_reauth),
        })


@app.get("/api/answer-quality")
async def get_answer_quality(session: str = "default", id: str = ""):
    """deferred answer_quality 결과 조회. id 지정 시 단건, 없으면 세션 전체."""
    try:
        from ai_engine.rag.quality_store import load_quality, get_quality
        if id:
            m = get_quality(session, id)
            return {"session": session, "id": id, "quality": m, "ready": m is not None}
        return {"session": session, "quality": load_quality(session)}
    except Exception as e:
        return {"session": session, "quality": None, "error": str(e)[:200]}


_AQ_TASKS = set()  # deferred answer_quality 태스크 참조 보관(GC 방지)


@app.post("/api/agents/run-stream")

async def run_agent_stream(request: Request):
    body = await request.json()
    prompt = body.get("prompt", "")
    model = body.get("model", "anthropic.claude-sonnet-4-5-20250929-v1:0")
    system_prompt = body.get("systemPrompt", "")
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
    project_path = body.get("projectPath", "")
    open_file = body.get("openFile", "")
    open_file_content = body.get("openFileContent", "")
    template_id = body.get("templateId", "")  # 활성 PPTX 템플릿 (요구사항 5.1)

    gw = _get_gw(aws_profile, bedrock_user)

    if project_path and not system_prompt:
        system_prompt = f"사용자의 프로젝트 경로: {project_path}"
        if open_file:
            system_prompt += f"\n현재 열린 파일: {open_file}"


    # 에이전트 모드 기본 지시 — 도구 사용 강제 + 작업 완료 후 보고
    _AGENT_BASE = (
        "\n\n[에이전트 모드 지시]\n"
        "- 사용자의 요청을 완수하기 위해 제공된 도구를 반드시 실행하세요.\n"
        "- 도구를 실행하지 않고 텍스트로만 답변하지 마세요.\n"
        "- 작업을 모두 완료한 후 결과를 간결하게 보고하세요.\n"
        "- 이미지/PDF/PPTX 생성 요청 시 generate_image/generate_pdf/generate_pptx 도구를 즉시 호출하세요.\n"
        "- 파일 읽기/쓰기/검색이 필요하면 read_file/write_file/search_files 도구를 사용하세요.\n"
    )
    system_prompt = (system_prompt or "") + _AGENT_BASE

    rag_evidence = None
    if project_path and _is_code_related(prompt):
        try:
            from ai_engine.rag.context_builder import build_system_prompt
            from ai_engine.rag.answer_quality import verify_mode as _vmode
            if _vmode(os.environ) != "off":
                system_prompt, rag_evidence = build_system_prompt(
                    project_path=project_path, query=prompt,
                    open_file=open_file, open_file_content=open_file_content,
                    base_system_prompt=system_prompt,
                    aws_profile=aws_profile, bedrock_user=bedrock_user, gateway_client=gw,
                    return_evidence=True,
                )
            else:
                system_prompt = build_system_prompt(
                    project_path=project_path, query=prompt,
                    open_file=open_file, open_file_content=open_file_content,
                    base_system_prompt=system_prompt,
                    aws_profile=aws_profile, bedrock_user=bedrock_user, gateway_client=gw,
                )
        except Exception as e:
            print(f"[RAG] 컨텍스트 빌드 실패 (무시): {e}")

    # 활성 템플릿 컨텍스트 주입 (요구사항 5.1) — 무템플릿이면 빈 문자열.
    _tpl_ctx = _active_template_prompt_context(template_id)
    if _tpl_ctx:
        system_prompt = (system_prompt or "") + _tpl_ctx

    messages = _build_messages(body.get("chatHistory", []), prompt, body.get("sessionId", "default"))
    stream_model = _resolve_callable_model_id(model, aws_profile, bedrock_user)

    async def realtime_stream():
        """Lambda SSE를 실시간으로 프론트엔드에 중계 — ChatGPT처럼 글자가 써지는 효과.
        max_tokens로 끊기면 자동으로 이어서 생성 (최대 5회)."""
        nonlocal messages
        max_continues = int(os.environ.get("AE_MAX_CONTINUES", "50"))
        # === OpenAI 모델 분기 (gateway-openai-models, 요구사항 5/6) ===
        # OpenAI Responses 라우트 모델은 토큰 스트리밍이 아닌 완성형 응답이므로,
        # route_openai_chat로 처리 후 텍스트를 단일 delta로 흘려보내 기존 SSE 소비
        # 코드와 호환시킨다. 비-OpenAI는 아래 기존 Bedrock 경로 그대로(바이트 동일).
        if is_openai_model(stream_model):
            _client_intent = (body.get("intent") or "").lower().strip()
            _txt = ""
            _agent_files = []
            try:
                # GPT가 스스로 도구(검색/읽기/생성)를 실행하도록 함수호출 루프 사용.
                _ares = await route_openai_agent(
                    gw, stream_model, messages, system_prompt=system_prompt,
                    project_path=project_path, aws_profile=aws_profile,
                    bedrock_user=bedrock_user, template_id=template_id,
                    max_iters=8, timeout=120,
                )
                _txt = _ares.get("text", "") or ""
                _agent_files = _ares.get("verified_files", []) or []
                if _txt:
                    yield f"data: {json.dumps({'text': _txt}, ensure_ascii=False)}\n\n"
            except Exception as _oe:
                # 도구 루프 실패 → 단발 응답으로라도 텍스트 시도(후 폴백).
                print(f"[run-stream/OpenAI] tool loop 실패 → 단발 폴백: {str(_oe)[:200]}")
                try:
                    _conv = await route_openai_chat(gw, stream_model, messages, system_prompt=system_prompt, timeout=120)
                    for _c in _conv.get("output", {}).get("message", {}).get("content", []):
                        if isinstance(_c, dict) and _c.get("text"):
                            _txt += _c["text"]
                    if _txt:
                        yield f"data: {json.dumps({'text': _txt}, ensure_ascii=False)}\n\n"
                except Exception as _oe2:
                    yield f"data: {json.dumps({'error': f'OpenAI 라우트 실패: {str(_oe2)[:200]}'}, ensure_ascii=False)}\n\n"
                    yield "data: [DONE]\n\n"
                    asyncio.create_task(_maybe_summarize(body.get("sessionId", "default"), body.get("chatHistory", []), gw))
                    return

            # 모델이 도구로 실제 파일을 만든 경우 → 즉시 통지(생성 탭 표시).
            _live_agent = []
            for _vf in _agent_files:
                _abs = _vf.get("absPath") or ""
                try:
                    if _abs and os.path.isfile(_abs) and os.path.getsize(_abs) > 0:
                        _live_agent.append(_vf.get("path", ""))
                except OSError:
                    continue
            if _live_agent:
                yield f"data: {json.dumps({'verifiedFiles': _live_agent}, ensure_ascii=False)}\n\n"

            # === 강제 생성 폴백 (gateway-openai-models) ===
            # 도구 루프로도 파일이 안 만들어졌고(또는 게이트웨이가 tool 미지원),
            # 사용자가 파일 생성을 요청한 정황이면 _force_generate_from_text로 복구.
            try:
                _pt, _wanted, _targets = _infer_file_intent_from_prompt(prompt, _client_intent, _txt)
                if _wanted and not _live_agent:
                    print(f"[run-stream/OpenAI] 도구 미생성 — 강제 생성 시도 (primary_tool={_pt})")
                    yield f"data: {json.dumps({'tool': 'deterministic-converter', 'status': 'running', 'input': {'reason': 'openai route — force generate'}}, ensure_ascii=False)}\n\n"
                    _force_start = __import__('time').time()
                    _enriched = ""
                    try:
                        _enrich_model = _specialized_model_for_task(
                            "enrich", stream_model, aws_profile=aws_profile, bedrock_user=bedrock_user)
                        _enriched = await _enrich_content_via_gateway(
                            gw=gw, model_id=_enrich_model, primary_tool=_pt,
                            title=prompt[:80], description=prompt,
                            final_text=_txt, project_path=project_path)
                    except Exception as _ee:
                        print(f"[run-stream/OpenAI] enrich 실패(무시): {_ee}")
                    if not _enriched or len(_enriched) < 200:
                        _real = _gather_real_context(prompt, project_path) \
                            or _gather_real_context_forced(project_path, prompt[:80])
                        if _real:
                            _enriched = _real
                    _forced = await _force_generate_from_text(
                        primary_tool=_pt, target_files=_targets,
                        title=prompt[:80], description=prompt,
                        final_text=_enriched or _txt or prompt,
                        project_path=project_path,
                        aws_profile=aws_profile, bedrock_user=bedrock_user,
                        template_id=template_id,
                    )
                    _force_ms = int((__import__('time').time() - _force_start) * 1000)
                    _live = []
                    for _fpath, _finfo in _forced:
                        _abs = _finfo.get("absPath") or ""
                        try:
                            if not _abs or not os.path.isfile(_abs) or os.path.getsize(_abs) <= 0:
                                continue
                        except OSError:
                            continue
                        _live.append(_finfo.get("path", ""))
                        yield f"data: {json.dumps({'tool': _finfo.get('tool', 'deterministic-converter'), 'path': _finfo.get('path', ''), 'output': _finfo.get('path', ''), 'status': 'done', 'durationMs': _force_ms}, ensure_ascii=False)}\n\n"
                        try:
                            _meta_obj = {
                                "tool": _finfo.get("tool", "deterministic"),
                                "model": _finfo.get("model", "deterministic-converter"),
                                "chatModel": stream_model,
                                "agentId": "single", "agentRole": "Agent",
                                "agentTitle": prompt[:80],
                                "createdAt": datetime.utcnow().isoformat() + "Z",
                                "promptHint": "[forced fallback/openai] " + (prompt or "")[:180],
                                "forced": True,
                            }
                            with open(_finfo["absPath"] + ".meta.json", "w", encoding="utf-8") as _mf:
                                json.dump(_meta_obj, _mf, ensure_ascii=False, indent=2)
                        except Exception:
                            pass
                    if _live:
                        _notice = (
                            "\n\n---\n실제 파일을 생성했습니다:\n" + "\n".join(f"- `{p}`" for p in _live)
                        )
                        yield f"data: {json.dumps({'text': _notice}, ensure_ascii=False)}\n\n"
                        yield f"data: {json.dumps({'verifiedFiles': _live}, ensure_ascii=False)}\n\n"
                    else:
                        _fail = (
                            "\n\n---\n⚠️ 파일 생성에 실패했습니다. 다시 시도하시거나 "
                            "\"PPTX로 저장해줘\"처럼 더 구체적으로 요청해 주세요."
                        )
                        yield f"data: {json.dumps({'text': _fail}, ensure_ascii=False)}\n\n"
            except Exception as _fe:
                print(f"[run-stream/OpenAI] 강제 생성 블록 예외: {_fe}")
            yield "data: [DONE]\n\n"
            asyncio.create_task(_maybe_summarize(body.get("sessionId", "default"), body.get("chatHistory", []), gw))
            return
        try:
            _full_answer_parts = []
            for cont in range(max_continues + 1):
                text_parts = []
                stop_reason = ""
                import time as _hb_time2
                _stream_start_ts = _hb_time2.time()
                def _mk_stream_s():
                    return gw.stream_sse_realtime(model_id=stream_model, messages=messages, system_prompt=system_prompt)
                async for evt in _stream_with_heartbeat(_mk_stream_s):
                    evt_type = evt.get("type", "")
                    if evt_type == "heartbeat":
                        _elapsed = int(_hb_time2.time() - _stream_start_ts)
                        yield f"data: {json.dumps({'heartbeat': True, 'elapsed': _elapsed, 'phase': 'thinking'}, ensure_ascii=False)}\n\n"
                        continue
                    if evt_type == "content_block_delta":
                        delta = evt.get("delta", {})
                        _rc = delta.get("reasoningContent") if isinstance(delta, dict) else None
                        if _rc and isinstance(_rc, dict):
                            _rtext = _rc.get("text", "")
                            if not _rtext and isinstance(_rc.get("reasoningText"), dict):
                                _rtext = _rc["reasoningText"].get("text", "")
                            if _rtext:
                                yield f"data: {json.dumps({'thinking': _rtext}, ensure_ascii=False)}\n\n"
                        if "text" in delta:
                            text_parts.append(delta["text"])
                            _full_answer_parts.append(delta["text"])
                            yield f"data: {json.dumps({'text': delta['text']}, ensure_ascii=False)}\n\n"
                    elif evt_type in ("message_delta", "message_stop"):
                        stop_reason = evt.get("delta", {}).get("stopReason", "") or evt.get("stop_reason", "") or evt.get("stopReason", "") or stop_reason
                    elif evt_type == "settlement":
                        rq = {"cost_krw": evt.get("remaining_quota_krw", 0)}
                        _extract_quota({"remaining_quota": rq, "estimated_cost_krw": evt.get("estimated_cost_krw", 0)}, _quota_cache.get("user", ""))
                    elif evt_type == "error":
                        yield f"data: {json.dumps({'error': evt.get('message', str(evt))}, ensure_ascii=False)}\n\n"

                # max_tokens로 끊김 → 이어서 생성
                if stop_reason == "max_tokens" and cont < max_continues and text_parts:
                    print(f"[Stream] max_tokens 도달 — 이어서 생성 ({cont+1}/{max_continues})")
                    messages.append({"role": "assistant", "content": [{"text": "".join(text_parts)}]})
                    messages.append({"role": "user", "content": [{"text": "계속 이어서 작성해주세요."}]})
                    continue
                break
        except Exception as e:
            yield f"data: {json.dumps({'error': str(e)}, ensure_ascii=False)}\n\n"
        # answer_quality (플래그 게이트). inline=최종이벤트로 대기, deferred=백그라운드+저장.
        try:
            from ai_engine.rag.answer_quality import (
                verify_mode as _vm, enhance_answer as _ea, run_deferred_verification as _rdv,
            )
            _mode = _vm(os.environ)
            _ans = "".join(_full_answer_parts).strip()
            if _mode != "off" and rag_evidence is not None and _ans:
                if _mode == "inline":
                    _r = await _ea(_ans, context_text=rag_evidence.get("context", ""),
                                   retrieved_chunks=rag_evidence.get("chunks"), gw=gw, env=os.environ)
                    _meta = _r.get("metadata") or {}
                    if _meta:
                        yield f"data: {json.dumps({'answerQuality': _meta}, ensure_ascii=False)}\n\n"
                elif _mode == "deferred":
                    import uuid as _uuid
                    _qid = _uuid.uuid4().hex
                    yield f"data: {json.dumps({'qualityPending': _qid}, ensure_ascii=False)}\n\n"
                    _t = asyncio.create_task(_rdv(_ans, rag_evidence.get("context", ""),
                                                  rag_evidence.get("chunks"), gw,
                                                  body.get("sessionId", "default"), _qid, os.environ))
                    _AQ_TASKS.add(_t); _t.add_done_callback(_AQ_TASKS.discard)
        except Exception as _aqe:
            print(f"[AnswerQuality] stream 검증 스킵(비차단): {_aqe}")

        yield "data: [DONE]\n\n"
        asyncio.create_task(_maybe_summarize(body.get("sessionId", "default"), body.get("chatHistory", []), gw))

    return StreamingResponse(
        realtime_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


# ─────────────────────────────────────────────────────────────────────────────
# LangGraph 그래프 경로 (Phase 1 골격) — feature flag `AE_LANGGRAPH`
#
# design.md 5단계 마이그레이션의 Phase 1. 기존 `/api/agents/run-stream` 과 **병행**
# 하는 신규 라우트 `/api/agents/graph-stream` 을 노출한다(요구사항 7.1).
#   - flag on:  coding 서브그래프를 compile 해 astream_events 로 실행(요구사항 7.2).
#   - flag off: 기존 run_agent_stream 핸들러로 위임(요구사항 7.3).
#   - 준비(compile/deps) 단계 실패: 기존 경로로 자동 fallback(요구사항 7.4).
#   - 스트리밍 시작 후 실패: `{error}` + `[DONE]` 로 종료(스트림은 되돌릴 수 없음).
#
# ⚠️ Phase 1 SSE 는 **최소 매핑**(on_chat_model_stream → {text}, 종료 → [DONE])만 한다.
#    완전한 이벤트 매핑(tool/verifiedFiles/agent_*/heartbeat)은 sse_bridge(Task 5.4)에서 교체.
# ⚠️ API_NOTES CRITICAL 2: 스트림 소비 루프(async for)를 asyncio.wait_for 로 감싸지 않는다
#    (Python 3.14 에서 취소 시 hang). per-node 타임아웃 + recursion_limit 로만 상한을 건다.
# ─────────────────────────────────────────────────────────────────────────────
def _langgraph_enabled() -> bool:
    """`AE_LANGGRAPH` 환경변수가 활성(1/true/on)인지 여부."""
    return os.environ.get("AE_LANGGRAPH", "").strip().lower() in ("1", "true", "on")


@app.post("/api/agents/graph-stream")
async def run_agent_graph_stream(request: Request):
    """LangGraph coding 서브그래프 기반 SSE 스트림 (Phase 1 골격).

    body 필드는 `/api/agents/run-stream` 과 동일(prompt/sessionId/projectPath/openFile/
    openFileContent/awsProfile/bedrockUser/templateId/model 등).
    """
    # flag 비활성 → 기존 경로로 위임(요구사항 7.3). body 재파싱은 FastAPI 가 캐시하므로 안전.
    if not _langgraph_enabled():
        return await run_agent_stream(request)

    body = await request.json()

    # ── 그래프 경로 준비(compile/deps). 이 단계 예외는 기존 경로로 fallback(요구사항 7.4). ──
    try:
        from langchain_core.messages import HumanMessage
        from ai_engine.agent_system.deps import GraphDeps
        from ai_engine.agent_system.supervisor import build_top_graph
        from ai_engine.agent_system.checkpoint_store import JsonFileCheckpointSaver
        from ai_engine.agent_system.sse_bridge import graph_events_to_sse

        prompt = body.get("prompt", "")
        model = body.get("model", "anthropic.claude-sonnet-4-5-20250929-v1:0")
        system_prompt = body.get("systemPrompt", "")
        aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
        bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
        project_path = body.get("projectPath", "")
        open_file = body.get("openFile", "")
        open_file_content = body.get("openFileContent", "")
        template_id = body.get("templateId", "")
        session_id = body.get("sessionId", "default")

        # GatewayClient — 기존 run-stream 과 동일한 방식(런타임 자격증명 주입, 파일 저장 없음).
        gw = _get_gw(aws_profile, bedrock_user)

        # checkpointer base_dir 는 userData 하위로 한정(요구사항 4.3 / 8.3).
        # Electron 이 주입한 AE_GENERATED_ROOT(userData) 하위 checkpoints/langgraph 를 우선.
        _env_root = os.environ.get("AE_GENERATED_ROOT", "").strip()
        if _env_root:
            _ckpt_dir = os.path.join(_env_root, "checkpoints", "langgraph")
        else:
            _ckpt_dir = ""  # JsonFileCheckpointSaver 기본(AE_CHECKPOINT_DIR / ~/.agentic-editor)
        checkpointer = JsonFileCheckpointSaver(_ckpt_dir)

        # Phase 2: Top Supervisor + graph-of-graphs(build_top_graph). deps 로 gateway/model/ckpt 주입.
        # (Phase 1 의 build_coding_subgraph 단일 그래프를 대체 — 요구사항 1.1/1.3/6.5.)
        deps = GraphDeps(
            gateway=gw,
            model_coding=_resolve_callable_model_id(model, aws_profile, bedrock_user),
            checkpointer=checkpointer,
        )
        compiled = build_top_graph(deps)

        # 초기 GraphState — 자격증명은 담지 않고 profile name / bedrock_user 문자열만(요구사항 8.1).
        # visited_routes 는 라우터 hop cap(요구사항 6.5)용 누적 리스트 — 명시적으로 [] 로 초기화한다
        # (reducer operator.add 기본값과 동일하나, 계약을 분명히 하기 위해 초기 상태에 포함).
        initial_state = {
            "prompt": prompt,
            "session_id": session_id,
            "project_path": project_path,
            "open_file": open_file,
            "open_file_content": open_file_content,
            "aws_profile": aws_profile,
            "bedrock_user": bedrock_user,
            "template_id": template_id,
            "system_prompt": system_prompt,
            "messages": [HumanMessage(content=prompt)],
            "visited_routes": [],
        }
        # config: thread_id(체크포인트 영속) + recursion_limit(요구사항 6.6, AE_GRAPH_RECURSION 기본 50).
        graph_config = {
            "configurable": {"thread_id": session_id},
            "recursion_limit": int(os.environ.get("AE_GRAPH_RECURSION", "50")),
        }
        # 그래프 전체 상한 / heartbeat 주기(요구사항 6.6/6.7/6.8, Property 4). ⚠️ API_NOTES
        # CRITICAL 2: 스트림 소비 루프를 asyncio.wait_for 로 감싸면 Python 3.14 에서 취소 시
        # hang → sse_bridge 가 개별 __anext__ 만 shield+wait 하고 deadline 을 수동 검사한다.
        # 각 상수는 이미 각 모듈에 env override 로 정의됨(중복 정의 금지, 재사용):
        #   - MAX_ROUTE_HOPS          : supervisor.py (AE_MAX_ROUTE_HOPS,      기본 4)
        #   - SUBGRAPH_RECURSION_LIMIT : subgraphs/_common.py (AE_SUBGRAPH_RECURSION, 기본 25)
        #   - MODEL_NODE_TIMEOUT       : subgraphs/_common.py (AE_MODEL_NODE_TIMEOUT)
        #   - recursion_limit          : 위 graph_config (AE_GRAPH_RECURSION,    기본 50)
        # 여기서는 그래프 전체 시간 상한 + heartbeat 주기를 sse_bridge 로 넘겨 배선한다.
        graph_total_timeout = float(os.environ.get("AE_GRAPH_TOTAL_TIMEOUT", "1800"))
        heartbeat_interval = float(os.environ.get("AE_HEARTBEAT_INTERVAL", "20"))
    except Exception as _prep_err:
        print(f"[graph-stream] 준비 실패 → 기존 run-stream 경로로 fallback: {_prep_err}")
        return await run_agent_stream(request)

    async def graph_stream():
        """Top Supervisor 그래프의 astream_events(v2) → sse_bridge 로 기존 SSE 계약 중계.

        SSE 매핑(on_chat_model_stream→{text}, tool→{tool,status}, 서브그래프 진입/종료→
        agent_start/agent_done, verifiedFiles, heartbeat, [DONE])과 노드 예외/
        GatewayModelError → {error}→[DONE] 처리, 그리고 무한대기 차단(개별 __anext__ 만
        shield+wait, deadline 수동 검사, 요구사항 5.7/6.6/6.7/6.8)은 모두 graph_events_to_sse
        내부가 담당한다. ⚠️ API_NOTES CRITICAL 2: 이 중계 루프를 asyncio.wait_for 로 감싸지
        않는다(스트림 제너레이터 취소 시 hang).
        """
        async for sse_line in graph_events_to_sse(
            compiled,
            initial_state,
            graph_config,
            heartbeat_interval=heartbeat_interval,
            total_timeout=graph_total_timeout,
        ):
            yield sse_line

    return StreamingResponse(
        graph_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


@app.post("/api/agents/run-agent")
async def run_agent_with_tools(request: Request):
    """에이전트 모드 — 도구 실행 루프 포함. 모델이 tool_use로 응답하면 실행 후 재호출.

    사용자가 도구 호출 미지원 모델(Llama/DeepSeek/Cohere 등)을 선택해도
    프롬프트가 도구 사용을 필요로 하면 Claude Sonnet으로 자동 라우팅한다.
    """
    body = await request.json()
    prompt = body.get("prompt", "")
    model = body.get("model", "anthropic.claude-sonnet-4-5-20250929-v1:0")
    system_prompt = body.get("systemPrompt", "")
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
    project_path = body.get("projectPath", "")
    open_file = body.get("openFile", "")
    open_file_content = body.get("openFileContent", "")
    template_id = body.get("templateId", "")  # 활성 PPTX 템플릿 (요구사항 5.1) — 도구 루프 _execute_tool에 전달

    # ── 자동 모델 라우팅 — Specialized Matrix ──
    # 사용자가 도구 호출 미지원/불안정 모델(Llama/DeepSeek/Cohere/Nova-Lite 등)을 선택했거나
    # 첨부 이미지가 있는데 비전 미지원 모델인 경우 자동으로 적절한 Claude로 대체.
    # Option A 원칙: 사용자 선택이 적합하면 그대로 유지, 부적합할 때만 매트릭스 적용.

    # chatHistory에 이미지 첨부가 있는지 검사 → 비전 capability 힌트
    def _detect_vision_in_history(history) -> bool:
        if not isinstance(history, list):
            return False
        for msg in history:
            if not isinstance(msg, dict):
                continue
            if msg.get("images"):
                return True
            content = msg.get("content")
            if isinstance(content, list):
                for blk in content:
                    if isinstance(blk, dict) and ("image" in blk or "imageSource" in blk):
                        return True
        return False

    has_vision = _detect_vision_in_history(body.get("chatHistory", []))
    caps_hint = ["vision"] if has_vision else None

    # 클라이언트가 intent classifier 결과를 전달하면 task 매핑에 사용 — 다양한 모델
    # 활용의 핵심. 비도구 작업(reasoning/translation/long_context/simple_qa)은
    # Claude 외 모델(DeepSeek/Qwen/Nova/Llama)로 라우팅된다.
    client_intent = (body.get("intent") or "").lower().strip()
    needs_tools_hint = bool(body.get("needs_tools"))
    # 프롬프트 자체로 파일/도구 생성 의도 추론 — intent 분류기 누락/오분류 대비.
    try:
        _pt_route, _wanted_route, _ = _infer_file_intent_from_prompt(prompt, client_intent, "")
    except Exception:
        _pt_route, _wanted_route = "", False
    if needs_tools_hint or _wanted_route:
        # 도구/파일 생성이 필요하면 반드시 도구 호출 가능 모델(Sonnet)로 라우팅한다.
        # 사용자가 NVIDIA/Gemma 등 도구 미지원 모델을 선택해도 빈 응답으로 실패하지 않도록
        # 특화(도구 가능) 모델로 자동 전환 → 생성이 항상 동작한다.
        task_for_routing = "file_generation"
    elif client_intent in _INTENT_TO_TASK:
        task_for_routing = _INTENT_TO_TASK[client_intent]
    else:
        task_for_routing = "general_chat"

    routed = _specialized_model_for_task(
        task_for_routing, model,
        capabilities_hint=caps_hint,
        aws_profile=aws_profile, bedrock_user=bedrock_user,
    )
    # 라우팅 결과가 사용자 선택과 다르면 알림 (raw id 비교)
    if routed and model and (model.lower().lstrip("us.").lstrip("eu.") not in routed.lower()):
        print(f"[Agent] 자동 라우팅 — {model} → {routed} (task={task_for_routing}, vision={has_vision})")
    model = routed

    gw = _get_gw(aws_profile, bedrock_user)
    stream_model = model  # 이미 prefix 처리된 callable id

    # 시스템 프롬프트 구성
    if project_path and not system_prompt:
        system_prompt = f"사용자의 프로젝트 경로: {project_path}"
        if open_file:
            system_prompt += f"\n현재 열린 파일: {open_file}"
    rag_evidence = None
    if project_path and _is_code_related(prompt):
        try:
            from ai_engine.rag.context_builder import build_system_prompt
            from ai_engine.rag.answer_quality import verify_mode as _vmode
            if _vmode(os.environ) != "off":
                system_prompt, rag_evidence = build_system_prompt(
                    project_path=project_path, query=prompt,
                    open_file=open_file, open_file_content=open_file_content,
                    base_system_prompt=system_prompt,
                    aws_profile=aws_profile, bedrock_user=bedrock_user, gateway_client=gw,
                    return_evidence=True,
                )
            else:
                system_prompt = build_system_prompt(
                    project_path=project_path, query=prompt,
                    open_file=open_file, open_file_content=open_file_content,
                    base_system_prompt=system_prompt,
                    aws_profile=aws_profile, bedrock_user=bedrock_user, gateway_client=gw,
                )
        except Exception as e:
            print(f"[Agent] RAG 실패 (무시): {e}")

    # 활성 템플릿 컨텍스트 주입 (요구사항 5.1) — 무템플릿이면 빈 문자열.
    _tpl_ctx = _active_template_prompt_context(body.get("templateId", ""))
    if _tpl_ctx:
        system_prompt = (system_prompt or "") + _tpl_ctx

    # 이슈 1 — 이미지 생성 의도면 generate_image 도구 사용을 명시적으로 강조.
    # 채팅 모델은 지휘자 역할이며, 도구가 내부에서 이미지 특화 모델을 자동 선택한다.
    try:
        _img_pt, _img_wanted, _ = _infer_file_intent_from_prompt(prompt, client_intent, "")
        if _img_pt == "generate_image":
            system_prompt = (system_prompt or "") + (
                "\n\n[이미지 생성 지시]\n"
                "- 사용자가 이미지 생성을 요청했습니다. 반드시 generate_image 도구를 호출하세요.\n"
                "- 텍스트로만 '생성했다'고 답하지 마세요. 도구를 실제로 실행해야 파일이 만들어집니다.\n"
                "- 이미지 모델은 시스템이 자동으로 최적 모델(Stable Image/Nova Canvas/Titan 등)을 선택합니다.\n"
            )
    except Exception:
        pass

    async def agent_stream():
        """에이전트 루프 — 최상위 try/finally 로 어떤 예외에도 [DONE] 송출 보장.
        ERR_INCOMPLETE_CHUNKED_ENCODING 방지 핵심."""
        # model_denied 재라우팅 시 enclosing 핸들러 스코프의 stream_model을 갱신한다.
        # (nonlocal 미선언 시 내부 재할당이 stream_model을 지역변수로 만들어 이전 읽기가
        #  UnboundLocalError를 일으킨다.)
        nonlocal stream_model
        done_sent = False
        # 자동 라우팅 알림 — body의 model과 stream_model이 다르면 사용자 알림
        _orig = body.get("model", "")
        if _orig and _orig.lower() not in (stream_model or "").lower():
            try:
                yield f"data: {json.dumps({'model_routing': True, 'original': _orig, 'routedTo': stream_model, 'reason': '도구 호출 안정성을 위해 Claude로 자동 라우팅됨'}, ensure_ascii=False)}\n\n"
            except Exception:
                pass
        try:
            # ── 메시지 구성 (실패해도 스트림은 이미 시작된 상태) ──
            try:
                messages = _build_messages(body.get("chatHistory", []), prompt, body.get("sessionId", "default"))
            except Exception as e:
                print(f"[Agent] _build_messages 실패: {e}")
                yield f"data: {json.dumps({'error': f'message build failed: {e}'}, ensure_ascii=False)}\n\n"
                return

            max_turns = int(os.environ.get("AE_MAX_AGENT_TURNS", "50"))
            refreshed_once = False  # 자격증명 만료 자동복구 1회만
            tool_unsupported_fallback_tried = False  # tool-use 미지원 모델 fallback 1회만
            model_denied_reroute_tried = False  # model_denied → Claude 재라우팅 1회만

            # === 단일 호출 파일 검증 추적 (이슈 2: 거짓 완료 방지) ===
            # 오케스트레이터(_orchestrator_run_agent_inner)와 동일한 원칙:
            # 모델이 도구를 실제로 호출해 디스크에 파일이 생겼는지 검증한다.
            # 도구를 안 부르고 텍스트로만 "생성 완료"라고 답하면 turn loop 종료 후
            # _force_generate_from_text로 실제 파일을 강제 생성한다.
            agent_verified_files = []   # [{path, absPath, size, tool}]
            agent_final_text_parts = []  # 모델이 누적한 본문 (강제 생성 시 enrich 입력)

            # === OpenAI 모델 분기 (gateway-openai-models, 요구사항 5/6) ===
            # OpenAI Responses 라우트는 Bedrock tool_config 루프와 호환되지 않으므로,
            # 완성형 응답으로 처리해 텍스트를 흘려보내고 turn loop를 건너뛴다.
            # 이후 강제 생성 폴백(파일 의도 시)은 기존 코드가 그대로 처리한다.
            # 비-OpenAI는 _openai_handled=False로 기존 turn loop 그대로(바이트 동일).
            _openai_handled = False
            if is_openai_model(stream_model):
                try:
                    # GPT가 스스로 도구(검색/읽기/생성)를 실행하는 함수호출 루프.
                    _ares = await route_openai_agent(
                        gw, stream_model, messages, system_prompt=system_prompt,
                        project_path=project_path, aws_profile=aws_profile,
                        bedrock_user=bedrock_user, template_id=body.get("templateId", ""),
                        max_iters=8, timeout=120,
                    )
                    _txt = _ares.get("text", "") or ""
                    if _txt:
                        agent_final_text_parts.append(_txt)
                        yield f"data: {json.dumps({'text': _txt}, ensure_ascii=False)}\n\n"
                    # 모델이 도구로 만든 파일을 검증 목록에 합류(후속 post-loop가 통지).
                    for _vf in (_ares.get("verified_files", []) or []):
                        _abs = _vf.get("absPath") or ""
                        try:
                            if _abs and os.path.isfile(_abs) and os.path.getsize(_abs) > 0 \
                                    and not any(v.get("absPath") == _abs for v in agent_verified_files):
                                agent_verified_files.append(_vf)
                                yield f"data: {json.dumps({'tool': _vf.get('tool', 'tool'), 'path': _vf.get('path', ''), 'output': _vf.get('path', ''), 'status': 'done'}, ensure_ascii=False)}\n\n"
                        except OSError:
                            continue
                except Exception as _oe:
                    print(f"[Agent/OpenAI] tool loop 실패 → 단발 폴백: {str(_oe)[:200]}")
                    try:
                        _conv = await route_openai_chat(gw, stream_model, messages, system_prompt=system_prompt, timeout=120)
                        _txt = ""
                        for _c in _conv.get("output", {}).get("message", {}).get("content", []):
                            if isinstance(_c, dict) and _c.get("text"):
                                _txt += _c["text"]
                        if _txt:
                            agent_final_text_parts.append(_txt)
                            yield f"data: {json.dumps({'text': _txt}, ensure_ascii=False)}\n\n"
                    except Exception as _oe2:
                        yield f"data: {json.dumps({'error': f'OpenAI 라우트 실패: {str(_oe2)[:200]}'}, ensure_ascii=False)}\n\n"
                _openai_handled = True

            for turn in range(0 if _openai_handled else max_turns):
                use_tool_config = not tool_unsupported_fallback_tried
                print(f"[Agent] turn={turn}, realtime stream, toolConfig={use_tool_config}")
                text_parts = []
                tool_use_blocks = []
                current_tool = {}
                stop_reason = ""
                turn_error = None
                need_reroute_retry = False  # model_denied 재라우팅 후 같은 turn 재시도 신호
                import time as _hb_time
                _turn_start_ts = _hb_time.time()

                try:
                    def _mk_stream():
                        return gw.stream_sse_realtime(
                            model_id=stream_model, messages=messages,
                            system_prompt=system_prompt,
                            tool_config=(AGENT_TOOLS if use_tool_config else None),
                        )
                    async for evt in _stream_with_heartbeat(_mk_stream):
                        evt_type = evt.get("type", "")
                        # heartbeat — 모델이 thinking 중이라 idle. 스트림을 살아있게 유지하고
                        # UI에 경과 시간을 표시(끊김 오판 방지). 텍스트로 누적하지 않는다.
                        if evt_type == "heartbeat":
                            _elapsed = int(_hb_time.time() - _turn_start_ts)
                            yield f"data: {json.dumps({'heartbeat': True, 'elapsed': _elapsed, 'phase': 'thinking'}, ensure_ascii=False)}\n\n"
                            continue
                        if evt_type == "content_block_delta":
                            delta = evt.get("delta", {})
                            # reasoning(thinking) 델타 — ChatGPT/Gemini처럼 생각 과정을 노출.
                            # 본문 텍스트로 누적하지 않고 별도 thinking 채널로만 전달한다.
                            _rc = delta.get("reasoningContent") if isinstance(delta, dict) else None
                            if _rc:
                                _rtext = ""
                                if isinstance(_rc, dict):
                                    _rtext = _rc.get("text") or (_rc.get("reasoningText") or {}).get("text", "") if isinstance(_rc.get("reasoningText"), dict) else _rc.get("text", "")
                                if _rtext:
                                    yield f"data: {json.dumps({'thinking': _rtext}, ensure_ascii=False)}\n\n"
                            if "text" in delta:
                                text_parts.append(delta["text"])
                                agent_final_text_parts.append(delta["text"])
                                yield f"data: {json.dumps({'text': delta['text']}, ensure_ascii=False)}\n\n"
                            elif "toolUse" in delta:
                                if current_tool:
                                    current_tool["_input_json"] = current_tool.get("_input_json", "") + delta["toolUse"].get("input", "")
                        elif evt_type == "content_block_start":
                            cb = evt.get("content_block") or evt.get("contentBlock") or {}
                            if "toolUse" in cb:
                                tu = cb["toolUse"]
                                current_tool = {"toolUseId": tu.get("toolUseId", ""), "name": tu.get("name", ""), "_input_json": ""}
                        elif evt_type == "content_block_stop":
                            if current_tool and current_tool.get("name"):
                                try:
                                    inp = json.loads(current_tool.get("_input_json", "{}"))
                                except json.JSONDecodeError:
                                    inp = {}
                                tool_use_blocks.append({
                                    "toolUse": {"toolUseId": current_tool["toolUseId"], "name": current_tool["name"], "input": inp}
                                })
                                current_tool = {}
                        elif evt_type in ("message_delta", "message_stop"):
                            stop_reason = evt.get("delta", {}).get("stopReason", "") or evt.get("stop_reason", "") or evt.get("stopReason", "")
                        elif evt_type == "settlement":
                            _extract_quota({"remaining_quota": {"cost_krw": evt.get("remaining_quota_krw", 0)}, "estimated_cost_krw": evt.get("estimated_cost_krw", 0)}, _quota_cache.get("user", ""))
                        elif evt_type == "error":
                            msg = evt.get("message", str(evt))
                            turn_error = msg
                            # 자격증명 만료 자동 복구 후 현재 turn 재시도
                            if (not refreshed_once) and _is_expired_error(msg):
                                refreshed_once = True
                                print("[Agent] 자격증명 만료 감지 — 자동 갱신 시도")
                                ok = await _refresh_and_retry_gw(gw, aws_profile, bedrock_user)
                                if ok:
                                    yield f"data: {json.dumps({'info': 'credentials refreshed, retrying'}, ensure_ascii=False)}\n\n"
                                    break  # async for 를 벗어나 현재 turn 재시도
                            # model_denied(미허용 모델) 자동 복구 — denylist 등록 후 Claude로 재라우팅.
                            _ml = str(msg).lower()
                            if ("model_denied" in _ml or "not in allowed list" in _ml) and not model_denied_reroute_tried:
                                model_denied_reroute_tried = True
                                _maybe_record_denied_from_error(msg)
                                _safe = _finalize_route_to_claude(stream_model, aws_profile, bedrock_user)
                                if _safe and _safe.lower() != (stream_model or "").lower():
                                    print(f"[Agent] model_denied 감지 — {stream_model} → {_safe} 재라우팅")
                                    stream_model = _safe
                                    need_reroute_retry = True
                                    yield f"data: {json.dumps({'model_routing': True, 'routedTo': _safe, 'reason': '미허용 모델 — Claude로 자동 재라우팅'}, ensure_ascii=False)}\n\n"
                                    break  # 현재 turn 을 안전 모델로 재시도
                            yield f"data: {json.dumps({'error': msg}, ensure_ascii=False)}\n\n"
                except asyncio.CancelledError:
                    # 클라이언트가 중단한 경우 — 조용히 종료
                    print("[Agent] 클라이언트 중단")
                    return
                except Exception as e:
                    import traceback
                    print(f"[Agent] stream 내부 예외: {e}\n{traceback.format_exc()}")
                    yield f"data: {json.dumps({'error': str(e)}, ensure_ascii=False)}\n\n"
                    break

                # 자격증명 갱신 후 재시도 케이스
                if turn_error and _is_expired_error(turn_error) and refreshed_once and not text_parts and not tool_use_blocks:
                    # 같은 turn 인덱스를 다시 쓰기 위해 range 를 못 돌리므로, messages 는 그대로 두고 continue 로 다음 turn 에서 재호출
                    continue

                # model_denied 재라우팅 후 같은 작업을 안전 모델(Claude)로 재시도
                if need_reroute_retry and not text_parts and not tool_use_blocks:
                    continue


                # tool-use 미지원 모델 fallback — InternalServerException + 아직 콘텐츠 없음 + 첫 시도
                if turn_error and not tool_unsupported_fallback_tried and not text_parts and not tool_use_blocks:
                    err_lower = str(turn_error).lower()
                    if "internalserverexception" in err_lower or "internal server" in err_lower or "tool" in err_lower:
                        tool_unsupported_fallback_tried = True
                        print(f"[Agent] tool-use 미지원 모델 감지 — toolConfig 없이 재시도")
                        yield f"data: {json.dumps({"info": "tool-unsupported, retrying without tools"}, ensure_ascii=False)}\n\n"
                        continue

                # content_blocks 조합
                content_blocks = []
                if text_parts:
                    content_blocks.append({"text": "".join(text_parts)})
                content_blocks.extend(tool_use_blocks)

                print(f"[Agent] turn={turn}, stopReason={stop_reason}, text={len(text_parts)}parts, tools={len(tool_use_blocks)}")

                if not content_blocks:
                    # Surface error or warn about empty response instead of silent break
                    if turn_error:
                        error_msg = f"모델 응답 오류: {str(turn_error)[:300]}"
                        print(f"[Agent] turn_error surfaced: {turn_error}")
                        yield f"data: " + json.dumps({"text": error_msg}, ensure_ascii=False) + "\n\n"
                    else:
                        print(f"[Agent] WARNING: empty response from model={stream_model}, sys_prompt_len={len(system_prompt)}")
                        yield f"data: " + json.dumps({"text": "⚠️ 모델이 빈 응답을 반환했습니다. 다시 시도해 주세요."}, ensure_ascii=False) + "\n\n"
                    break
                messages.append({"role": "assistant", "content": content_blocks})

                if not tool_use_blocks:
                    if stop_reason == "max_tokens" and turn < max_turns - 1:
                        print(f"[Agent] max_tokens 도달 — 이어서 생성 (turn {turn+1})")
                        messages.append({"role": "user", "content": [{"text": "계속 이어서 작성해주세요."}]})
                        continue
                    break

                # 도구 실행
                tool_results = []
                for block in tool_use_blocks:
                    tu = block["toolUse"]
                    tool_name = tu.get("name", "")
                    tool_id = tu.get("toolUseId", "")
                    tool_input = tu.get("input", {})
                    yield f"data: {json.dumps({'tool': tool_name, 'input': tool_input, 'status': 'running'}, ensure_ascii=False)}\n\n"
                    import time as _time
                    _tool_start = _time.time()
                    # 도구 실행을 태스크로 돌리고, 길어지면 heartbeat를 주기적으로 송출한다.
                    # (이미지/PPTX 생성 등은 수 분 걸릴 수 있어 idle 워치독 오판을 막아야 함)
                    _tool_task = asyncio.ensure_future(
                        asyncio.to_thread(_execute_tool, tool_name, tool_input, project_path, aws_profile, bedrock_user, template_id)  # [patched-credentials]
                    )
                    try:
                        while True:
                            try:
                                tool_output = await asyncio.wait_for(asyncio.shield(_tool_task), timeout=_SSE_HEARTBEAT_SECONDS)
                                break
                            except asyncio.TimeoutError:
                                _elapsed = int(_time.time() - _tool_start)
                                yield f"data: {json.dumps({'heartbeat': True, 'elapsed': _elapsed, 'phase': 'tool', 'tool': tool_name}, ensure_ascii=False)}\n\n"
                                continue
                    except Exception as e:
                        tool_output = f"도구 실행 예외: {e}"
                    _tool_duration_ms = int((_time.time() - _tool_start) * 1000)
                    print(f"[Agent] 도구 실행: {tool_name} → {len(tool_output)}자 ({_tool_duration_ms}ms)")

                    # 생성 파일에 .meta.json 사이드카 기록
                    if tool_name in ("generate_image", "generate_pdf", "generate_pptx", "generate_xlsx", "generate_docx", "edit_image", "write_file"):
                        try:
                            _meta_paths = []
                            _actual_model = ""
                            try:
                                _parsed = json.loads(tool_output) if isinstance(tool_output, str) else None
                            except (json.JSONDecodeError, TypeError):
                                _parsed = None
                            if isinstance(_parsed, dict) and "path" in _parsed and "error" not in _parsed:
                                _meta_paths.append(_parsed["path"])
                                _actual_model = _parsed.get("model") or ""
                            elif isinstance(_parsed, dict) and isinstance(_parsed.get("images"), list):
                                for _it in _parsed["images"]:
                                    if isinstance(_it, dict) and "path" in _it:
                                        _meta_paths.append(_it["path"])
                                        if not _actual_model and _it.get("model"):
                                            _actual_model = _it["model"]
                            if tool_name == "write_file" and "path" in tool_input:
                                _meta_paths.append(tool_input["path"])
                            _tool_default_model = {
                                "generate_pdf":  "reportlab (Python)",
                                "generate_pptx": "python-pptx",
                                "generate_xlsx": "openpyxl",
                                "generate_docx": "python-docx",
                                "write_file":    "filesystem",
                            }
                            for _rel in _meta_paths:
                                _abs = _resolve_relative_for_verify(_rel, project_path)
                                if not os.path.isfile(_abs):
                                    continue
                                # 이슈 2 — 실제 디스크 존재 확인된 파일만 검증 목록에 추가.
                                # 이 목록이 비어있으면 turn loop 종료 후 강제 생성이 발동한다.
                                try:
                                    _vsize = os.path.getsize(_abs)
                                except OSError:
                                    _vsize = 0
                                if not any(vf.get("absPath") == _abs for vf in agent_verified_files):
                                    agent_verified_files.append({
                                        "path": _rel, "absPath": _abs,
                                        "size": _vsize, "tool": tool_name,
                                    })
                                _model_label = _actual_model or _tool_default_model.get(tool_name) or stream_model
                                _meta_obj = {
                                    "tool": tool_name,
                                    "model": _model_label,
                                    "chatModel": stream_model,
                                    "agentId": "single",
                                    "agentRole": "Agent",
                                    "agentTitle": prompt[:80],
                                    "createdAt": datetime.utcnow().isoformat() + "Z",
                                    "promptHint": prompt[:200],
                                }
                                try:
                                    with open(_abs + ".meta.json", "w", encoding="utf-8") as _mf:
                                        json.dump(_meta_obj, _mf, ensure_ascii=False, indent=2)
                                except Exception as _me:
                                    print(f"[Agent] meta write 실패 {_abs}: {_me}")
                        except Exception as _e:
                            print(f"[Agent] meta 처리 예외: {_e}")

                    yield f"data: {json.dumps({'tool': tool_name, 'output': tool_output[:500], 'status': 'done', 'durationMs': _tool_duration_ms}, ensure_ascii=False)}\n\n"
                    _tr_max = int(os.environ.get("AE_TOOL_RESULT_MAX", "80000"))
                    tool_results.append({"toolResult": {"toolUseId": tool_id, "content": [{"text": tool_output[:_tr_max]}]}})

                messages.append({"role": "user", "content": tool_results})

            # === 이슈 2: 도구 미호출 거짓 완료 방지 — 강제 생성 폴백 ===
            # turn loop가 끝났는데 디스크에 검증된 파일이 0개이고, 사용자가
            # 파일 생성을 요청한 정황이면 → 오케스트레이터와 동일하게
            # 게이트웨이 본문 보강 후 _force_generate_from_text로 실제 파일 생성.
            # 모델이 "생성했다"고 텍스트로만 답하고 실제로는 도구를 안 부른 케이스를 복구.
            try:
                _final_text = "".join(agent_final_text_parts)
                _pt, _wanted, _targets = _infer_file_intent_from_prompt(
                    prompt, client_intent, _final_text
                )
                if _wanted and not agent_verified_files:
                    print(f"[Agent] 도구 미호출 감지 — 강제 생성 시도 (primary_tool={_pt})")
                    yield f"data: {json.dumps({'tool': 'deterministic-converter', 'status': 'running', 'input': {'reason': 'no tool calls — force generate'}}, ensure_ascii=False)}\n\n"
                    _force_start = __import__('time').time()
                    try:
                        _enrich_model = _specialized_model_for_task(
                            "enrich", stream_model,
                            aws_profile=aws_profile, bedrock_user=bedrock_user,
                        )
                        _enriched = await _enrich_content_via_gateway(
                            gw=gw, model_id=_enrich_model, primary_tool=_pt,
                            title=prompt[:80], description=prompt,
                            final_text=_final_text, project_path=project_path,
                        )
                        if not _enriched or len(_enriched) < 200:
                            _real = _gather_real_context(prompt, project_path) \
                                or _gather_real_context_forced(project_path, prompt[:80])
                            if _real:
                                _enriched = _real
                        _forced = await _force_generate_from_text(
                            primary_tool=_pt, target_files=_targets,
                            title=prompt[:80], description=prompt,
                            final_text=_enriched or _final_text or prompt,
                            project_path=project_path,
                            aws_profile=aws_profile, bedrock_user=bedrock_user,
                            template_id=body.get("templateId", ""),
                        )
                        _force_ms = int((__import__('time').time() - _force_start) * 1000)
                        # TASK 8 — forced 결과를 한 번 더 디스크 재검증.
                        # _force_generate_from_text 내부도 isfile 체크하지만,
                        # 응답~프론트 emit 사이의 race(cleanup hook, antivirus 등)로
                        # 파일이 사라질 수 있어 emit 직전에 재확인한다.
                        for _fpath, _finfo in _forced:
                            _abs_check = _finfo.get("absPath") or ""
                            try:
                                if not _abs_check or not os.path.isfile(_abs_check) or os.path.getsize(_abs_check) <= 0:
                                    print(f"[Agent] forced 결과 재검증 실패 — 스킵: {_abs_check}")
                                    continue
                            except OSError:
                                continue
                            agent_verified_files.append(_finfo)
                            # meta sidecar — forced 경로도 실제 라이브러리 이름 유지
                            try:
                                _meta_obj = {
                                    "tool": _finfo.get("tool", "deterministic"),
                                    "model": _finfo.get("model", "deterministic-converter"),
                                    "chatModel": stream_model,
                                    "agentId": "single", "agentRole": "Agent",
                                    "agentTitle": prompt[:80],
                                    "createdAt": datetime.utcnow().isoformat() + "Z",
                                    "promptHint": "[forced fallback] " + (prompt or "")[:180],
                                    "forced": True,
                                }
                                with open(_finfo["absPath"] + ".meta.json", "w", encoding="utf-8") as _mf:
                                    json.dump(_meta_obj, _mf, ensure_ascii=False, indent=2)
                            except Exception:
                                pass
                            yield f"data: {json.dumps({'tool': _finfo.get('tool', 'deterministic-converter'), 'path': _finfo.get('path', ''), 'output': _finfo.get('path', ''), 'status': 'done', 'durationMs': _force_ms}, ensure_ascii=False)}\n\n"
                        if agent_verified_files:
                            _paths = [vf["path"] for vf in agent_verified_files]
                            _notice = (
                                "\n\n---\n실제 파일을 생성했습니다 (모델이 도구를 호출하지 않아 "
                                "시스템이 자동 생성):\n" + "\n".join(f"- `{p}`" for p in _paths)
                            )
                            yield f"data: {json.dumps({'text': _notice}, ensure_ascii=False)}\n\n"
                            yield f"data: {json.dumps({'verifiedFiles': _paths}, ensure_ascii=False)}\n\n"
                        else:
                            # TASK 8 — forced도 실패했으면 "생성됐다"는 모델 답변과 실제가
                            # 어긋난다. 사용자가 다운로드해도 파일 없음. 채팅에 명시.
                            print("[Agent] 강제 생성 결과 0건 — 사용자에게 실패 안내")
                            _fail_notice = (
                                "\n\n---\n⚠️ 파일 생성에 실패했습니다.\n"
                                "모델이 \"생성 완료\"라고 답했지만 실제로 디스크에 저장된 파일이 없습니다 "
                                "(도구 호출이 누락됐고 강제 생성 폴백도 실패).\n"
                                "다시 시도하시거나, 더 구체적으로 \"PPTX로 저장해줘\" 같이 요청해 주세요."
                            )
                            yield f"data: {json.dumps({'text': _fail_notice}, ensure_ascii=False)}\n\n"
                    except Exception as _fe:
                        print(f"[Agent] 강제 생성 실패: {_fe}")
                        yield f"data: {json.dumps({'tool': 'deterministic-converter', 'status': 'done', 'output': f'force-generate failed: {_fe}'}, ensure_ascii=False)}\n\n"
                elif agent_verified_files:
                    # 정상 경로로 파일이 생성된 경우에도 검증 목록을 프론트에 통지.
                    # TASK 8 — emit 직전 디스크 자체 재검증으로 사라진 파일 제외.
                    _live_paths = []
                    for vf in agent_verified_files:
                        _abs = vf.get("absPath") or _resolve_relative_for_verify(vf.get("path", ""), project_path)
                        try:
                            if _abs and os.path.isfile(_abs) and os.path.getsize(_abs) > 0:
                                _live_paths.append(vf["path"])
                        except OSError:
                            continue
                    if _live_paths:
                        yield f"data: {json.dumps({'verifiedFiles': _live_paths}, ensure_ascii=False)}\n\n"
                    else:
                        # 도구 응답엔 path가 있었는데 emit 시점에 모두 사라짐 — 명확히 알림.
                        print("[Agent] 도구 호출 후 디스크 재검증 0건")
                        _fail_notice = (
                            "\n\n---\n⚠️ 파일이 디스크에 남아있지 않습니다.\n"
                            "도구가 path를 응답했으나 실제 파일을 찾을 수 없습니다. "
                            "다시 시도해 주세요."
                        )
                        yield f"data: {json.dumps({'text': _fail_notice}, ensure_ascii=False)}\n\n"
            except Exception as _ve:
                print(f"[Agent] 파일 검증/강제생성 블록 예외: {_ve}")
            # answer_quality (플래그 게이트, inline/deferred, additive·비차단)
            try:
                from ai_engine.rag.answer_quality import (
                    verify_mode as _vm, enhance_answer as _ea, run_deferred_verification as _rdv,
                )
                _mode = _vm(os.environ)
                _ans = "".join(agent_final_text_parts).strip()
                if _mode != "off" and rag_evidence is not None and _ans:
                    if _mode == "inline":
                        _r = await _ea(_ans, context_text=rag_evidence.get("context", ""),
                                       retrieved_chunks=rag_evidence.get("chunks"), gw=gw, env=os.environ)
                        _meta = _r.get("metadata") or {}
                        if _meta:
                            yield f"data: {json.dumps({'answerQuality': _meta}, ensure_ascii=False)}\n\n"
                    elif _mode == "deferred":
                        import uuid as _uuid
                        _qid = _uuid.uuid4().hex
                        yield f"data: {json.dumps({'qualityPending': _qid}, ensure_ascii=False)}\n\n"
                        _t = asyncio.create_task(_rdv(_ans, rag_evidence.get("context", ""),
                                                      rag_evidence.get("chunks"), gw,
                                                      body.get("sessionId", "default"), _qid, os.environ))
                        _AQ_TASKS.add(_t); _t.add_done_callback(_AQ_TASKS.discard)
            except Exception as _aqe:
                print(f"[AnswerQuality] agent 검증 스킵(비차단): {_aqe}")
        except asyncio.CancelledError:
            print("[Agent] stream cancelled")
            return
        except Exception as e:
            import traceback
            print(f"[Agent] 최상위 예외: {e}\n{traceback.format_exc()}")
            try:
                yield f"data: {json.dumps({'error': f'fatal: {e}'}, ensure_ascii=False)}\n\n"
            except Exception:
                pass
        finally:
            # 어떤 경로로 끝나든 [DONE] 을 반드시 송출 → ERR_INCOMPLETE_CHUNKED_ENCODING 방지
            if not done_sent:
                done_sent = True
                try:
                    yield "data: [DONE]\n\n"
                except Exception:
                    pass

    return StreamingResponse(
        agent_stream(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


@app.post("/api/agents/run-parallel")
async def run_agent_parallel(request: Request):
    """병렬 모델 호출 — 서버에서 동시 실행, SSE로 각 모델 결과 전달."""
    body = await request.json()
    prompt = body.get("prompt", "")
    _raw_models = body.get("models", [])
    # models 정규화 — 항목이 문자열(modelId)로 와도 dict로 변환.
    # (디버그 print의 s.get('modelId')가 문자열에서 크래시 → SSE 스트림 전체가
    #  죽어 모든 슬롯이 에러로 표시되던 문제 방지. 프론트/구버전 무관하게 견고.)
    models = []
    for _i, _mm in enumerate(_raw_models if isinstance(_raw_models, list) else []):
        if isinstance(_mm, str):
            models.append({"modelId": _mm, "slotId": f"slot-{_i+1}", "systemPrompt": ""})
        elif isinstance(_mm, dict):
            models.append({
                "modelId": _mm.get("modelId") or _mm.get("model") or "",
                "slotId": _mm.get("slotId") or f"slot-{_i+1}",
                "systemPrompt": _mm.get("systemPrompt", "") or "",
                **{k: v for k, v in _mm.items() if k not in ("modelId", "slotId", "systemPrompt")},
            })
        # 그 외 타입은 건너뜀(방어)
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
    project_path = body.get("projectPath", "")
    open_file = body.get("openFile", "")
    open_file_content = body.get("openFileContent", "")
    template_id = body.get("templateId", "")  # 활성 PPTX 템플릿 (요구사항 5.1) — 도구 루프 _execute_tool에 전달

    gw = _get_gw(aws_profile, bedrock_user)

    # RAG 컨텍스트 — 코드/프로젝트 관련 질문에만
    rag_context = ""
    if project_path and _is_code_related(prompt):
        try:
            from ai_engine.rag.context_builder import build_system_prompt
            import functools as _ft
            _rag_fn = _ft.partial(
                build_system_prompt,
                project_path=project_path,
                query=prompt,
                open_file=open_file,
                open_file_content=open_file_content,
                aws_profile=aws_profile,
                bedrock_user=bedrock_user,
                gateway_client=gw,
            )
            # 대형 프로젝트에서 첫 RAG 인덱싱/임베딩은 수십 초~분이 걸린다.
            # 동기 호출은 이벤트 루프를 블로킹해 SSE 응답이 지연되고 모든 슬롯이
            # 타임아웃(에러)된다. 스레드에서 실행 + 예산 초과 시 RAG 없이 진행
            # (인덱스는 백그라운드에서 계속 구축 → 다음 호출부터 캐시 적중).
            rag_context = await asyncio.wait_for(
                asyncio.get_event_loop().run_in_executor(None, _rag_fn),
                timeout=float(os.environ.get("AE_RAG_BUDGET_S", "12")),
            )
        except asyncio.TimeoutError:
            print("[RAG] 컨텍스트 빌드 예산 초과 — RAG 없이 진행 (백그라운드 인덱싱 지속, 다음 호출 캐시)")
        except Exception as e:
            print(f"[RAG] 컨텍스트 빌드 실패 (무시): {e}")

    # 이전 대화 맥락(chatHistory) + 세션 메모리를 반영
    chat_history = body.get("chatHistory", [])
    session_id = body.get("sessionId", "default")
    messages = _build_messages(chat_history, prompt, session_id)

    async def parallel_stream():
        # 할루시네이션 sanitizer — 도구 호출 시뮬레이션과 거짓 완료 주장을 제거
        def _sanitize_hallucination(text: str) -> tuple:
            """모델이 도구 사용을 시뮬레이션한 흔적이나 거짓 주장을 감지하고 정리.
            실제 파일 존재 여부를 .generated/ 폴더에서 검증하여 거짓 경로는 제거.
            반환: (cleaned_text, was_modified)
            """
            if not text:
                return text, False
            original = text

            # 0) 실제 .generated/ 폴더에 어떤 파일이 있는지 확인
            _generated_dir = os.path.join(project_path or os.getcwd(), ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
            _existing_files = set()
            try:
                if os.path.isdir(_generated_dir):
                    for f in os.listdir(_generated_dir):
                        _existing_files.add(f)
            except Exception:
                pass

            # 1) function_calls/tool_call XML 태그 제거 (앞뒤 공백 포함)
            text = re.sub(r'<function_calls>.*?</function_calls>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<invoke[^>]*>.*?</invoke>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<tool_call>.*?</tool_call>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<parameter[^>]*>.*?</parameter>', '', text, flags=re.DOTALL | re.IGNORECASE)
            text = re.sub(r'<function_calls>|</function_calls>|<invoke[^>]*>|</invoke>|<tool_call>|</tool_call>', '', text, flags=re.IGNORECASE)

            # 2) 가짜 .generated/ 경로 검증 — 실제 존재하지 않는 파일은 표시 변경
            def _replace_fake_path(match):
                full = match.group(0)
                fname = match.group(1)
                if fname in _existing_files:
                    return full  # 실제 존재 → 그대로
                return f"`(없음: {fname})`"  # 가짜 → 명시적으로 없다고 표시
            text = re.sub(r'`?\.generated/([\w\-.]+)`?', _replace_fake_path, text)

            # 3) ![이미지](.generated/xxx.png) 마크다운 — 실제 파일 검증
            def _replace_md_img(match):
                fname = match.group(1)
                if fname in _existing_files:
                    return match.group(0)
                return f"[*이미지 없음: {fname}*]"
            text = re.sub(r'!\[[^\]]*\]\(\.generated/([\w\-.]+)\)', _replace_md_img, text)

            # 4) 표 형태의 거짓 완료 주장 제거 — "✅ 완료" 행
            # | 1 | PNG 이미지 | .generated/xxx.png | ✅ 완료 |  같은 패턴
            text = re.sub(r'\|[^\n|]*?(✅|완료|완성|생성됨|created|done)[^\n|]*?\|', '| (실제 파일 없음) |', text)

            # 5) 거짓 완료 주장 약화
            text = re.sub(r'(✅|✓|🎉)\s*[^\n.]{0,30}(생성|저장|작성|만들).{0,10}(완료|되었|됨)[!.]?', '*다음 내용으로 생성 가능:*', text, flags=re.IGNORECASE)
            text = re.sub(r'(이전 작업|위 작업).*?(이미|모두)\s*완료[^.]*\.?', '*(주의: 이전 대화에서 실제 파일이 생성되지 않았습니다.)*', text)

            # 6) "모든 파일이 .generated/에 저장됩니다" 같은 거짓 안내
            text = re.sub(r'(모든\s*)?파일.{0,20}\.generated/.{0,30}(저장|생성|만들|작성).{0,20}(됩니다|되었|완료|돼요)', '*아래 코드를 실행하면 .generated/에 파일이 생성됩니다*', text, flags=re.IGNORECASE)

            modified = (text != original)
            # 7) 너무 많이 잘려나갔으면 경고 추가
            if modified and len(text.strip()) < 50 and len(original) > 200:
                text = ("*[알림] 이 모델의 응답에서 도구 호출 시뮬레이션 또는 거짓 주장이 감지되어 정리되었습니다. "
                        "실제로는 파일이 생성되지 않았습니다. "
                        "실제 파일 생성을 원하시면 '에이전트로 작업 진행' 버튼을 사용하세요.*\n\n"
                        + text)
            return text.strip(), modified

        async def call_model(slot):
            model_id = slot.get("modelId", "")
            slot_id = slot.get("slotId", "")
            sp = slot.get("systemPrompt", "")
            # CRIS profile 존재 여부에 따라 us. prefix 적용 (ON_DEMAND 모델은 prefix 없이)
            sid = _resolve_callable_model_id(model_id, aws_profile, bedrock_user)
            # RAG 컨텍스트를 시스템 프롬프트에 추가
            if rag_context:
                sp = (sp + "\n\n" + rag_context) if sp else rag_context

            # 할루시네이션 방지 — 거짓 주장은 막되, 실제 내용은 충실히 작성하도록 균형
            _anti = (
                "\n\n[병렬 모드 응답 가이드 — 매우 중요]\n"
                "이 호출에서는 파일 생성/수정 도구를 사용할 수 없습니다.\n"
                "이전 대화에 '파일이 생성되었다' 같은 응답이 있었더라도 **그것은 거짓**이었습니다 — 실제로 .generated/ 폴더에 파일은 만들어지지 않았습니다.\n"
                "\n"
                "올바른 응답 방식:\n"
                "1. 사용자가 원하는 결과물의 **실제 내용**(코드, 텍스트, 구조, 설계)을 마크다운으로 직접 작성\n"
                "2. 코드 블록(```python, ```markdown 등)으로 충분히 상세하게 제공\n"
                "3. 사용자가 그대로 활용 가능한 완성도 높은 내용\n"
                "\n"
                "절대 금지 사항:\n"
                "- '파일을 생성했습니다', '저장 완료', '✅ 완료' 같은 거짓 주장\n"
                "- '.generated/xxx.png' 같은 가짜 파일 경로 출력\n"
                "- 표 형태로 '✅ 완료', '생성됨' 같은 가짜 상태 표시\n"
                "- <function_calls>, <invoke>, <tool_call> 등 도구 호출 시뮬레이션\n"
                "- '이전 작업이 모두 완료되었습니다' 같은 과거 거짓 사실 인정\n"
                "\n"
                "사용자가 원하는 것은 **실제 내용**입니다. '도구가 없다'는 변명만 하지 말고 내용을 충실히 작성하세요.\n"
                "사용자는 다음 단계로 '에이전트 모드'를 사용해 실제 파일을 만들 수 있으며, 당신의 응답이 그 원천이 됩니다."
            )
            sp = (sp + _anti) if sp else _anti.strip()

            # 병렬 실행 검증 로그
            import time as _time
            _t_start = _time.time()
            print(f"[Parallel] START slot={slot_id} model={model_id} t=0.000s")

            # === OpenAI 모델 분기 (gateway-openai-models) ===
            # 병렬 모드는 '도구 없이 내용 비교'가 원칙이므로 OpenAI도 단발 텍스트 응답.
            # Bedrock converse 경로(아래)는 OpenAI 모델을 처리하지 못하므로 분기 필수.
            if is_openai_model(sid):
                try:
                    _conv = await route_openai_chat(gw, sid, messages, system_prompt=sp, timeout=120)
                    _txt = ""
                    for _c in _conv.get("output", {}).get("message", {}).get("content", []):
                        if isinstance(_c, dict) and _c.get("text"):
                            _txt += _c["text"]
                    _txt, _ws = _sanitize_hallucination(_txt)
                    print(f"[Parallel] DONE(OpenAI) slot={slot_id} model={model_id} elapsed={_time.time()-_t_start:.2f}s")
                    return {"slotId": slot_id, "modelId": model_id, "status": "done", "content": _txt}
                except Exception as _oe:
                    return {"slotId": slot_id, "modelId": model_id, "status": "error",
                            "content": f"OpenAI 라우트 실패: {str(_oe)[:200]}"}

            # 자동 재시도 (최대 3회, 지수 백오프)
            for attempt in range(3):
                try:
                    result = await asyncio.wait_for(
                        gw.converse_stream_live(model_id=sid, messages=messages, system_prompt=sp),
                        timeout=600
                    )
                    if result.get("decision") == "ERROR":
                        err = result.get("error", "")
                        # throttling/rate limit → 재시도
                        if attempt < 2 and ("throttl" in err.lower() or "rate" in err.lower() or "timed out" in err.lower()):
                            await asyncio.sleep(2 ** attempt * 2)
                            continue
                        result = await asyncio.wait_for(
                            gw.converse(model_id=sid, messages=messages, system_prompt=sp),
                            timeout=600
                        )
                    decision = result.get("decision", "")
                    if decision == "ALLOW":
                        output = result.get("output", {}).get("message", {}).get("content", [])
                        text = "\n".join(c.get("text", "") for c in output if "text" in c)
                        # 할루시네이션 후처리 — 도구 호출 시뮬레이션, 거짓 완료 주장 제거
                        text, _was_sanitized = _sanitize_hallucination(text)
                        if _was_sanitized:
                            print(f"[Parallel] SANITIZED slot={slot_id} model={model_id}")
                        print(f"[Parallel] DONE  slot={slot_id} model={model_id} elapsed={_time.time()-_t_start:.2f}s")
                        return {"slotId": slot_id, "modelId": model_id, "status": "done", "content": text}
                    elif decision == "ACCEPTED":
                        job_id = result.get("job_id", "")
                        if job_id:
                            text = await gw._poll_job_result(job_id, max_wait=600)
                            if text:
                                return {"slotId": slot_id, "modelId": model_id, "status": "done", "content": text}
                        return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": "ACCEPTED — 결과 대기 시간 초과"}
                    elif decision == "DENY":
                        return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": result.get("denial_reason", "DENIED")}
                    else:
                        err_msg = result.get("error", f"Unknown: {decision}")
                        if attempt < 2 and ("throttl" in err_msg.lower() or "timed out" in err_msg.lower()):
                            await asyncio.sleep(2 ** attempt * 2)
                            continue
                        return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": err_msg}
                except asyncio.TimeoutError:
                    if attempt < 2:
                        await asyncio.sleep(2 ** attempt * 2)
                        continue
                    return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": "600초 타임아웃"}
                except Exception as e:
                    if attempt < 2:
                        await asyncio.sleep(2 ** attempt)
                        continue
                    return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": str(e)}
            return {"slotId": slot_id, "modelId": model_id, "status": "error", "content": "재시도 3회 실패"}

        # 배치 실행: 10개씩 동시 호출 (rate limit 방지) + 배치 간 heartbeat
        # asyncio.create_task로 명시적 즉시 스케줄링 → 진짜 동시 실행 보장
        import time as _time
        batch_size = 10
        _collected = []  # 병렬 결과 수집 — 요약 체크포인트 반영용
        for i in range(0, len(models), batch_size):
            batch = models[i:i+batch_size]
            _batch_t0 = _time.time()
            print(f"[Parallel] BATCH start size={len(batch)} models={[(s.get('modelId') if isinstance(s, dict) else s) for s in batch]}")
            # 모든 task를 즉시 스케줄링 — 이벤트 루프가 다음 await에서 모두 시작
            tasks = [asyncio.create_task(call_model(slot)) for slot in batch]
            # as_completed로 완료 순서대로 yield
            for coro in asyncio.as_completed(tasks):
                result = await coro
                _collected.append(result)
                yield f"data: {json.dumps(result, ensure_ascii=False)}\n\n"
            print(f"[Parallel] BATCH done elapsed={_time.time()-_batch_t0:.2f}s")
            # 배치 간 heartbeat — 클라이언트 idle timeout 방지
            if i + batch_size < len(models):
                yield f"data: {json.dumps({'heartbeat': True, 'progress': min(i+batch_size, len(models)), 'total': len(models)})}\n\n"

        # === 병렬 연속성 — 결과를 세션 요약 체크포인트에 반영 ===
        # 병렬 assistant 응답(합본)을 chat_history에 이어붙여 요약을 트리거한다.
        # 다음 턴의 _build_messages가 이 요약 체크포인트를 주입하므로, 병렬만 계속
        # 이어가는 긴 세션에서도 초반 맥락이 밀려나지 않고 유지된다(요구사항: 체크포인트
        # 컨텍스트 유지). 비동기 task로 던져 SSE 종료를 막지 않는다.
        try:
            _ok = [r for r in _collected
                   if isinstance(r, dict) and r.get("status") == "done" and r.get("content")]
            if _ok:
                _digest_parts = []
                for r in _ok:
                    _nm = r.get("modelId", "model")
                    _digest_parts.append(f"[{_nm}]\n{(r.get('content') or '')[:800]}")
                _digest = "\n\n".join(_digest_parts)[:4000]
                _summ_hist = list(chat_history or []) + [
                    {"role": "user", "content": prompt},
                    {"role": "assistant", "content": _digest},
                ]
                asyncio.create_task(_maybe_summarize(session_id, _summ_hist, gw))
        except Exception as _se:
            print(f"[Parallel] 요약 체크포인트 트리거 실패(무시): {_se}")

        yield "data: [DONE]\n\n"

    return StreamingResponse(parallel_stream(), media_type="text/event-stream")


# ─────────────────────────────────────────────────────────────────
# Multi-Agent Orchestrator  (Role Split + Parallel Tool-Using Agents)
# ─────────────────────────────────────────────────────────────────

ORCHESTRATOR_PLANNER_PROMPT = """당신은 멀티-에이전트 시스템의 플래너입니다.
사용자의 요청을 서로 독립적으로 병렬 실행 가능한 N개의 하위 작업(subtask)으로 분해하세요.

규칙:
1. 각 subtask는 서로 파일/코드 영역이 겹치지 않아야 합니다 (충돌 방지).
2. 각 subtask는 하나의 에이전트가 도구(read_file, write_file, list_directory, run_command, search_files, generate_image, generate_pdf, generate_pptx, generate_xlsx, generate_docx)를 사용해 독립적으로 완료할 수 있어야 합니다.
3. subtask 개수는 요청 내용에 맞게 결정하되, 최대 {max_agents}개를 넘지 마세요.
4. 사용자가 "수정 1~4" 처럼 명시한 번호/단계가 있으면 그대로 따르세요.
5. **파일 생성 요청의 경우 — 절대 규칙**:
   - 사용자가 명시한 **모든 파일 형식마다 별도 subtask**를 만드세요.
   - "PDF, PPTX, DOCX, XLSX 만들어줘" → 정확히 4개 subtask (PDF 1개 + PPTX 1개 + DOCX 1개 + XLSX 1개)
   - 형식을 임의로 통합/축소/누락하면 안 됩니다 (PNG/SVG로 대체 금지).
   - 형식별 매핑은 아래 6번 primary_tool과 정확히 1:1로 일치해야 합니다.
   - **사용자가 명시하지 않은 형식은 절대 추가하지 마세요**. 특히 PNG/이미지는 사용자가 "이미지", "PNG", "그림", "image", "그려줘", "diagram" 같은 단어를 명시했을 때만 추가 가능. PDF/PPTX/DOCX 만들어달라고만 했으면 PNG subtask 추가 금지.
6. **primary_tool ↔ 파일 형식 매핑 (엄격)**:
   - PDF (.pdf)   → primary_tool: "generate_pdf"
   - PPTX (.pptx) → primary_tool: "generate_pptx"
   - XLSX (.xlsx) → primary_tool: "generate_xlsx"
   - DOCX (.docx) → primary_tool: "generate_docx"
   - PNG/JPG (.png/.jpg) → primary_tool: "generate_image"
   - 이미지 편집 → primary_tool: "edit_image"
   - 코드/마크다운/텍스트 (.md/.py/.js/.txt) → primary_tool: "write_file"
   - 셸 명령 → primary_tool: "run_command"
   - 코드 분석/리팩토링 → primary_tool: "code"
7. 각 subtask의 description에는 **반드시 도구를 사용해 실제 파일을 생성/저장**하라고 명시하세요. "텍스트로 출력"이 아니라 "<primary_tool> 도구로 .generated/ 폴더에 저장".
8. target_files에는 결과물 파일의 절대/상대 경로를 명시하세요. **확장자는 primary_tool과 일치해야 함** (예: primary_tool=generate_xlsx → ".generated/...xlsx").

[출력 규칙 — 매우 중요]
- 응답은 **순수 JSON 객체 1개만**. 그 외 어떤 텍스트도 추가 금지.
- 마크다운 코드블록(```json, ```) 사용 금지.
- "여기 JSON입니다:" 같은 설명문 금지.
- 응답 첫 글자는 `{{`, 마지막 글자는 `}}`이어야 함.
- 여러 JSON 객체 출력 금지 — 정확히 1개만.

JSON 형식:
{{
  "subtasks": [
    {{
      "id": "A",
      "role": "역할명 (예: PDF Generator, XLSX Builder, Image Creator, Code Refactorer)",
      "title": "간결한 제목",
      "description": "이 에이전트가 수행해야 할 작업의 상세 지시. 반드시 도구를 사용해 실제 파일을 생성/저장하라고 명시.",
      "primary_tool": "generate_image|generate_pdf|generate_pptx|generate_xlsx|generate_docx|edit_image|write_file|run_command|code",
      "team": "research|coding|writing|media|analysis (선택, 같은 team 작업끼리 컨텍스트 공유)",
      "target_files": [".generated/result.pdf", ...]
    }},
    ...
  ]
}}
"""

ORCHESTRATOR_AGENT_PROMPT = """당신은 멀티-에이전트 시스템의 전문 작업자입니다.

[할당된 역할] {role}
[작업 ID] {task_id}
[작업 제목] {title}
[대상 파일] {target_files}
[핵심 도구] {primary_tool}

[지시사항]
{description}

[작업 규칙 — 매우 중요]
1. **반드시 도구를 사용하세요**. 텍스트로만 답변하지 마세요. 다음 도구들을 적극 활용:
   - read_file: 기존 파일 내용 확인
   - list_directory: 폴더 구조 파악
   - search_files: 코드/문서 검색
   - run_command: 셸 명령 실행 (Python 스크립트 등)
   - write_file: 텍스트/코드 파일 작성
   - generate_image: 이미지 생성 — 시스템이 프롬프트를 분석해 최적 모델을 자동 선택
     · 사진/리얼리즘 → Stability Stable Image Ultra
     · 다이어그램/차트/UI → Stability SD 3.5 Large (텍스트 렌더링 우수)
     · 일러스트/아트 → Stability Ultra
     · 로고/아이콘 → Stable Image Core (빠름)
     · 비즈니스 자료 → Amazon Nova Canvas
     **프롬프트는 영어로 작성하면 모델 자동 선택 정확도가 향상됩니다.**
   - generate_pdf:  PDF 생성 (reportlab)
   - generate_pptx: PPTX 생성 (python-pptx)
   - generate_xlsx: XLSX 엑셀 생성 (openpyxl) — sheets[{{name, headers, rows}}] 형태로 호출
   - generate_docx: DOCX 워드 생성 (python-docx) — sections[{{heading, body, bullets}}] 형태로 호출
   - edit_image:    이미지 inpaint/outpaint

2. **파일 형식과 도구는 정확히 일치해야 합니다 (절대 규칙)**:
   - .pdf  → generate_pdf  (절대 generate_image 금지)
   - .pptx → generate_pptx (절대 generate_image 금지)
   - .xlsx → generate_xlsx (절대 generate_image 금지)
   - .docx → generate_docx (절대 generate_image 금지)
   - .png/.jpg → generate_image
   사용자가 "PDF + PPTX + XLSX + DOCX 만들어줘"라고 했으면 위 4개 도구를 각각 호출하세요. PNG/SVG로 대체하면 작업 실패입니다.

3. **절대 금지 — 우회 시도 감지**:
   - primary_tool=generate_pdf 인데 write_file로 .py / .html / .md 작성 → 작업 실패로 처리됩니다.
   - primary_tool=generate_pptx 인데 generate_image로 PNG/SVG 만들기 → 작업 실패로 처리됩니다.
   - "PDF는 reportlab으로 만들어야 하니 코드를 작성합니다" → 금지. generate_pdf 도구가 이미 reportlab을 사용합니다. 도구를 직접 호출하세요.
   - "다이어그램은 SVG가 좋으니 SVG로 만듭니다" → 금지. 사용자가 PDF/PPTX를 요청했으면 그 형식으로만 만드세요.

4. **위 [핵심 도구]가 명시되어 있다면 가장 먼저 그 도구부터 호출하세요**. 텍스트 설명은 도구 호출 후에 추가하세요.

5. **이미지 생성 작업의 경우**: 사용자가 어떤 채팅 모델(Claude 등)을 선택했더라도, 실제 이미지는 generate_image 도구가 시스템 내부에서 Stability/Amazon 이미지 모델을 자동 호출해서 생성합니다. Claude로 이미지를 직접 그리려고 하지 마세요 — generate_image 도구를 호출하면 됩니다.

6. **파일 생성 작업의 표준 절차**:
   a. 필요시 list_directory/read_file로 컨텍스트 수집
   b. 적절한 generate_* 또는 write_file 도구로 실제 파일 생성
   c. 생성된 파일 경로를 응답에 포함
   d. **반드시 .generated/ 폴더에 저장** (없으면 run_command로 mkdir)

7. **대상 파일 외의 파일은 수정하지 마세요**.

8. **다른 에이전트의 작업 영역을 침범하지 마세요**.

9. 작업이 끝나면 "[완료] <생성된 파일 경로 + 한 줄 요약>" 형태로 마무리하세요.

10. 도구 호출 없이 텍스트만 출력하면 작업이 실패한 것으로 간주됩니다. 반드시 도구를 사용해 실제 결과물을 만들어내세요.
"""

ORCHESTRATOR_MERGER_PROMPT = """당신은 멀티-에이전트 결과를 통합하는 리뷰어입니다.

[입력 데이터의 verifiedFiles 필드는 시스템이 디스크에서 직접 확인한 실제 파일 목록입니다 — 이 목록만 신뢰하세요.]
[에이전트의 summary 텍스트에 "✅ 완료" "생성됨" 등이 적혀 있어도 verifiedFiles에 없으면 실패입니다.]

아래 각 에이전트의 작업 결과를 검토하고:
1. 성공/실패 여부 — verifiedFiles에 파일이 있고 status="done"이면 성공, 그 외는 모두 실패
2. 충돌이나 누락 지적
3. 사용자에게 전달할 최종 요약 보고서를 마크다운으로 작성

보고서 형식:
## 최종 통합 결과
| 에이전트 | 역할 | 상태 | 생성된 파일 (디스크 검증됨) | 도구 사용 횟수 |
|---------|------|------|--------------------------|--------------|
...

### 생성된 파일 목록 (verifiedFiles 기반 — 실제 존재)
- `.generated/file1.pdf` — 설명
- `.generated/file2.xlsx` — 설명
...

### 세부 사항
- (각 에이전트별 핵심 작업 내역)

### 주의/후속 작업
- (있다면)

**규칙**:
- verifiedFileCount=0인 에이전트는 무조건 "실패"로 표시.
- summary에 "✅", "완료", "생성됨"이 있어도 verifiedFiles에 없으면 그 주장은 무시하고 실패로 분류.
- 거짓 파일 경로(verifiedFiles에 없는 경로)는 절대 보고서에 포함하지 마세요.

**오류 추정 절대 금지**:
- 입력 데이터에 명시되지 않은 오류 메시지를 만들어내지 마세요.
- 'KeyError', 'TypeError', 'ValidationException' 같은 구체적 예외 이름은 입력 summary에 그대로 적혀있을 때만 인용하세요.
- summary에 오류가 안 적혀 있으면 "도구를 호출하지 않음" 또는 "원인 불명"으로만 기재하세요.
- 절대 그럴듯해 보이는 가짜 KeyError나 스택트레이스를 작성하지 마세요.

**사용자 명시 형식 검증**:
- summary_input.requiredFormats 필드가 있으면, verifiedFiles 의 확장자가 그 형식들을 모두 포함하는지 확인하세요.
- 누락된 형식이 있으면 보고서에 "**누락**: PDF, PPTX 등" 형태로 명시하고 "전체 성공"이라 쓰지 마세요.
"""


def _extract_first_json_object(text: str) -> dict:
    """LLM 응답에서 첫 번째 valid JSON 객체를 안전하게 추출.

    처리하는 케이스:
    1. 마크다운 코드블록 (```json ... ``` 또는 ``` ... ```)
    2. 응답 앞뒤의 자유 텍스트 ("여기 JSON입니다:" 등)
    3. 여러 JSON 객체가 줄바꿈으로 나란히 있는 경우
    4. 중첩 중괄호가 있는 큰 JSON 객체
    5. 문자열 안에 있는 중괄호 (escape 처리)
    6. 'subtasks' 키가 있는 객체 우선 선택

    Returns:
        파싱된 dict. 실패 시 빈 dict {} 반환.
    """
    if not text:
        return {}

    # Step 1: 마크다운 코드블록 제거
    # ```json\n{...}\n``` 또는 ```\n{...}\n``` 패턴
    code_block_match = re.search(r"```(?:json)?\s*\n?([\s\S]*?)\n?```", text)
    if code_block_match:
        text = code_block_match.group(1).strip()

    # Step 2: 모든 top-level JSON 객체 추출 (brace counting + string-aware)
    candidates = []
    n = len(text)
    i = 0
    while i < n:
        if text[i] == "{":
            # JSON 객체 시작 — 매칭되는 } 찾기
            depth = 0
            in_string = False
            escape = False
            start = i
            j = i
            while j < n:
                ch = text[j]
                if escape:
                    escape = False
                elif ch == "\\":
                    escape = True
                elif in_string:
                    if ch == '"':
                        in_string = False
                elif ch == '"':
                    in_string = True
                elif ch == "{":
                    depth += 1
                elif ch == "}":
                    depth -= 1
                    if depth == 0:
                        # 완전한 JSON 객체 후보
                        candidate = text[start:j + 1]
                        try:
                            parsed = json.loads(candidate)
                            if isinstance(parsed, dict):
                                candidates.append(parsed)
                        except json.JSONDecodeError:
                            pass
                        i = j + 1
                        break
                j += 1
            else:
                # 닫는 } 못 찾음 — break
                break
        else:
            i += 1

    if not candidates:
        return {}

    # Step 3: 'subtasks' 키가 있는 객체 우선
    for c in candidates:
        if "subtasks" in c and isinstance(c["subtasks"], list):
            return c
    # 없으면 첫 번째 객체
    return candidates[0]


# ─────────────────────────────────────────────────────────────────
# 사용자 요청 형식 추출 + Planner/파이프라인 보강 헬퍼
# ─────────────────────────────────────────────────────────────────

# 형식 → 도구 매핑 (실제 generator가 있는 형식만 enforcement 대상)
_REQUIRED_FORMAT_TOOL = {
    "pdf": "generate_pdf",
    "pptx": "generate_pptx",
    "xlsx": "generate_xlsx",
    "docx": "generate_docx",
    "png": "generate_image",
}


def _fmt_hits_in_segment(seg_text: str, seg_raw: str) -> set:
    """한 세그먼트(문장/구)에서 언급된 파일 형식 키워드를 추출.

    seg_text = seg_raw.lower(). 실제 generator 도구가 있는 형식만 반환:
    pdf, pptx, xlsx, docx, png. (SVG/MD는 코드 유사 출력이라 제외.)
    """
    hits = set()
    if re.search(r"(?:^|[^a-z])pdf(?:[^a-z]|$)", seg_text) or ".pdf" in seg_text:
        hits.add("pdf")
    if (re.search(r"(?:^|[^a-z])pptx(?:[^a-z]|$)", seg_text) or ".pptx" in seg_text
            or "파워포인트" in seg_raw or "프레젠테이션" in seg_raw
            or "슬라이드" in seg_raw or "powerpoint" in seg_text):
        hits.add("pptx")
    if (re.search(r"(?:^|[^a-z])xlsx(?:[^a-z]|$)", seg_text) or ".xlsx" in seg_text
            or "엑셀" in seg_raw or "스프레드시트" in seg_raw or "excel" in seg_text):
        hits.add("xlsx")
    if (re.search(r"(?:^|[^a-z])docx(?:[^a-z]|$)", seg_text) or ".docx" in seg_text
            or "워드" in seg_raw or re.search(r"(?:^|[^a-z])word(?:[^a-z]|$)", seg_text)):
        hits.add("docx")
    if (re.search(r"(?:^|[^a-z])(?:png|jpg|jpeg|image)(?:[^a-z]|$)", seg_text)
            or "이미지" in seg_raw or "그림" in seg_raw or "사진" in seg_raw):
        hits.add("png")
    return hits


# 파일 생성을 명시적으로 거부/제외하는 부정 문맥 신호.
# 이 신호가 형식 키워드와 같은 세그먼트에 있으면 그 형식은 "요청"으로 보지 않는다.
_FMT_NEG_RE = re.compile(
    r"(하지\s*마|하지마|마세요|마셈|만들지\s*(?:마|말|않)|생성\s*(?:하지|치\s*마|금지|안\s)"
    r"|필요\s*없|필요없|말고|제외|금지|없이|빼고|말아|불필요"
    r"|do\s*not|don['’]?t|without|no\s+need|not\s+(?:create|make|generate|need|required))"
)
# 실제 파일 생성을 요청하는 긍정 의도 신호. 형식 키워드가 "요청"으로 인정되려면
# 같은 세그먼트에 이 신호가 함께 있어야 한다(단순 언급/부정 맥락 오탐 방지).
_FMT_POS_RE = re.compile(
    r"(생성|만들|작성|제작|그려|그리|출력|저장|뽑아|제출|내보|변환|첨부"
    r"|해\s*줘|해줘|해\s*주|부탁|원해|원합|주세요|줘"
    r"|create|make|generate|draw|write|build|export|produce|convert|want|need|give\s+me)"
)


def _extract_required_formats(user_prompt: str) -> set:
    """사용자 프롬프트에서 *명시적으로 생성 요청된* 파일 형식들을 추출.

    강제성 완화(사용자 요청): 형식 키워드가 프롬프트에 등장하기만 해서는
    강제 생성 대상으로 삼지 않는다. 세그먼트(문장/구) 단위로 판정하며,
    다음을 모두 만족하는 형식만 반환한다:
      1) 같은 세그먼트에 파일 생성을 요청하는 *긍정 의도*(만들/생성/작성 등)가 있고,
      2) 같은 세그먼트에 *부정/거부 문맥*(하지 마/없이/제외/금지 등)이 없다.

    이로써 "요청하지 않은 파일(PDF/PPTX 등) 생성은 하지 마세요" 같은 부정문이나,
    이전 답변/로그가 섞인 종합 프롬프트의 단순 형식 언급이 강제 생성을 유발하지
    않는다. 실제 generator 도구가 있는 형식(pdf/pptx/xlsx/docx/png)만 대상.

    Returns:
        set[str] — lowercase 확장자 집합. 명시 요청이 없으면 빈 set.
    """
    if not user_prompt:
        return set()

    # 문장/구 단위 분할 — 형식과 부정/긍정 신호의 지역성(locality)을 판정하기 위함.
    # 주의: 긍정 동사(만들고 등)를 분할자로 쓰면 앞 세그먼트가 긍정 신호를 잃으므로
    # 문장부호와 명확한 접속사만 분할자로 사용한다. 복합 대조문
    # ("A는 만들고 B는 만들지 마")은 부정 우선으로 보수적 처리(강제 생성 안 함) —
    # planner가 문맥으로 보완하므로 안전하다.
    segments = re.split(r"[.!?\n,;·:•]|그리고|또한|하지만|반면", user_prompt)
    requested = set()
    excluded = set()
    for seg_raw in segments:
        seg_text = seg_raw.lower()
        hits = _fmt_hits_in_segment(seg_text, seg_raw)
        if not hits:
            continue
        if _FMT_NEG_RE.search(seg_raw) or _FMT_NEG_RE.search(seg_text):
            # 부정 문맥 → 이 형식은 강제하지 않음(오탐 방지)
            excluded |= hits
        elif _FMT_POS_RE.search(seg_raw) or _FMT_POS_RE.search(seg_text):
            # 긍정 생성 의도가 명확할 때만 요청으로 인정
            requested |= hits
        # 긍정도 부정도 없는 단순 언급은 강제하지 않음(보수적 — 강제성 제거)

    return requested - excluded


def _subtask_primary_format(st: dict) -> str:
    """subtask에서 결과물 1차 파일 형식을 추출. 형식이 명확하지 않으면 ''."""
    if not isinstance(st, dict):
        return ""
    pt = (st.get("primary_tool") or "").lower()
    tool_to_fmt = {
        "generate_pdf": "pdf",
        "generate_pptx": "pptx",
        "generate_xlsx": "xlsx",
        "generate_docx": "docx",
        "generate_image": "png",
    }
    if pt in tool_to_fmt:
        return tool_to_fmt[pt]
    # target_files 첫 항목의 확장자
    tf = st.get("target_files") or []
    if isinstance(tf, list) and tf:
        first = str(tf[0]).lower()
        ext = first.rsplit(".", 1)[-1] if "." in first else ""
        if ext in ("pdf", "pptx", "xlsx", "docx", "png"):
            return ext
        if ext in ("jpg", "jpeg"):
            return "png"
    return ""


def _subtask_is_neutral(st: dict) -> bool:
    """subtask가 파일 산출 없이 연구/분석만 하는 중립 작업인지 판정."""
    if not isinstance(st, dict):
        return False
    pt = (st.get("primary_tool") or "").lower()
    return pt in {"read_file", "list_directory", "search_files", "run_command"}


def _subtask_format_matches(st: dict, required: set) -> bool:
    """subtask의 1차 형식이 사용자 요청 형식 집합 안에 있는지."""
    fmt = _subtask_primary_format(st)
    return bool(fmt) and fmt in (required or set())


def _synthetic_subtask(fmt: str, user_prompt: str) -> dict:
    """누락 형식을 보강하기 위한 강제 subtask 생성. role/title은 한국어."""
    fmt = (fmt or "").lower()
    tool = _REQUIRED_FORMAT_TOOL.get(fmt, "")
    label_map = {
        "pdf": ("PDF 보고서 작성자", "PDF 강제 생성"),
        "pptx": ("PPTX 슬라이드 작성자", "PPTX 강제 생성"),
        "xlsx": ("XLSX 워크북 작성자", "XLSX 강제 생성"),
        "docx": ("DOCX 문서 작성자", "DOCX 강제 생성"),
        "png": ("이미지 생성자", "PNG 이미지 강제 생성"),
    }
    role, title = label_map.get(fmt, (f"{fmt.upper()} 작성자", f"{fmt.upper()} 강제 생성"))
    # 안전한 파일명 슬러그 — 사용자 프롬프트의 앞부분에서 ASCII만 추려 사용
    safe = re.sub(r"[^A-Za-z0-9._-]+", "-", (user_prompt or "result").strip())[:40] or "result"
    return {
        "id": f"FORCE-{fmt.upper()}",
        "role": role,
        "title": title,
        "description": (
            f"사용자가 명시적으로 요청한 {fmt.upper()} 형식이 다른 subtask에서 누락되었습니다. "
            f"반드시 {tool} 도구를 호출해서 .generated/ 폴더에 {fmt.upper()} 파일을 생성하세요. "
            f"내용은 사용자 원본 요청을 기반으로 작성: {user_prompt[:400]}"
        ),
        "primary_tool": tool,
        "team": "writing" if fmt in ("pdf", "pptx", "xlsx", "docx") else "media",
        "target_files": [f".generated/forced-{safe}.{fmt}"],
    }


async def _orchestrator_plan(gw, stream_model, user_prompt: str, system_prompt: str, max_agents: int) -> dict:
    """Planner 호출 — subtask 분해. 응답 JSON 파싱 실패에 대한 자동 재시도 포함."""
    plan_sys = ORCHESTRATOR_PLANNER_PROMPT.format(max_agents=max_agents)
    if system_prompt:
        plan_sys = system_prompt + "\n\n" + plan_sys

    last_text = ""
    last_err = ""

    for attempt in range(3):
        # 첫 시도 외엔 추가 지시문으로 강화
        if attempt == 0:
            messages = [{"role": "user", "content": [{"text": user_prompt}]}]
        else:
            retry_instruction = (
                f"이전 응답에서 JSON을 파싱할 수 없었습니다 ({last_err[:120]}).\n"
                f"이번에는 반드시 다음 규칙을 지켜 주세요:\n"
                f"1. 마크다운 코드블록(```) 사용 금지\n"
                f"2. JSON 앞뒤에 어떤 텍스트도 추가 금지\n"
                f"3. 단 하나의 JSON 객체만 출력\n"
                f"4. {{\"subtasks\": [...]}} 형태\n\n"
                f"원본 요청: {user_prompt}"
            )
            messages = [{"role": "user", "content": [{"text": retry_instruction}]}]

        try:
            result = await asyncio.wait_for(
                gw.converse(model_id=stream_model, messages=messages, system_prompt=plan_sys),
                timeout=120,
            )
        except Exception as e:
            last_err = f"converse 예외: {e}"
            continue

        if result.get("decision") != "ALLOW":
            last_err = str(result.get("error") or result.get("decision"))
            continue

        output = result.get("output", {}).get("message", {}).get("content", [])
        text = "\n".join(c.get("text", "") for c in output if "text" in c).strip()
        last_text = text

        # 안전한 JSON 추출
        plan = _extract_first_json_object(text)
        if plan and isinstance(plan.get("subtasks"), list) and plan["subtasks"]:
            print(f"[Planner] {attempt+1}회 시도 성공 — {len(plan['subtasks'])}개 subtask")
            return plan

        last_err = "subtasks 키가 없거나 비어있음" if plan else "JSON 객체 찾을 수 없음"
        print(f"[Planner] {attempt+1}회 시도 실패 — {last_err}, 응답 일부: {text[:200]}")

    # 3회 모두 실패 — 사용자 요청을 1개 subtask로 wrap (절대 실패하지 않음)
    print(f"[Planner] 3회 시도 모두 실패. 단일 subtask fallback. last_text={last_text[:300]}")
    return {
        "subtasks": [{
            "id": "A",
            "role": "General Worker",
            "title": "사용자 요청 수행",
            "description": user_prompt[:500],
            "primary_tool": "",
            "target_files": [],
        }],
    }


def _gather_real_context(description: str, project_path: str) -> str:
    """description에서 작업 유형을 감지하여 실제 데이터를 수집.

    LLM이 추측하지 않도록 소스에서 직접 데이터를 모은다.
    여러 카테고리의 데이터를 동시에 수집해 합성 가능 (예: "폴더 구조 + README").

    감지 카테고리:
    1. folder_structure  - 폴더 트리 + 통계
    2. dependencies      - package.json/requirements.txt/Cargo.toml 등
    3. git_summary       - 최근 커밋 + 변경 파일
    4. readme_summary    - README/CHANGELOG 본문
    5. code_inventory    - 언어별 파일 개수 + 주요 파일 (LOC 상위)
    6. config_summary    - 주요 설정 파일 (.gitignore, dockerfile 등)

    Returns:
        합성된 마크다운 텍스트. 모든 카테고리 빈 결과면 빈 문자열.
    """
    if not project_path or not os.path.isdir(project_path):
        return ""
    desc_lower = (description or "").lower()
    sections = []

    # ─── 1. Folder structure ───
    folder_keywords = (
        "폴더 구조", "디렉토리 구조", "폴더 트리", "디렉토리 트리",
        "폴더 깊이", "뎁스", "depth", "흐름도", "계층",
        "folder structure", "directory structure", "tree",
        "프로젝트 구조", "project structure", "파일 트리",
    )
    if any(kw in desc_lower for kw in folder_keywords):
        s = _gather_folder_tree(project_path)
        if s:
            sections.append(s)

    # ─── 2. Dependencies ───
    dep_keywords = (
        "의존성", "의존 라이브러리", "dependency", "dependencies",
        "package", "requirements", "라이브러리", "외부 라이브러리",
        "버전", "스택", "기술 스택", "tech stack",
    )
    if any(kw in desc_lower for kw in dep_keywords):
        s = _gather_dependencies(project_path)
        if s:
            sections.append(s)

    # ─── 3. Git summary ───
    git_keywords = (
        "git", "커밋", "commit", "변경 이력", "히스토리",
        "최근 변경", "changelog", "변경 사항", "history",
    )
    if any(kw in desc_lower for kw in git_keywords):
        s = _gather_git_summary(project_path)
        if s:
            sections.append(s)

    # ─── 4. README/CHANGELOG ───
    readme_keywords = (
        "readme", "리드미", "프로젝트 소개", "프로젝트 개요",
        "프로젝트 설명", "overview", "introduction", "소개",
        "changelog", "변경 로그", "릴리즈 노트", "release note",
    )
    if any(kw in desc_lower for kw in readme_keywords):
        s = _gather_readme(project_path)
        if s:
            sections.append(s)

    # ─── 5. Code inventory ───
    code_keywords = (
        "코드 분석", "code analysis", "loc", "라인 수",
        "코드베이스", "codebase", "주요 파일", "main file",
        "큰 파일", "largest file", "코드 통계", "코드 통계",
        "분석", "analyze", "리뷰", "review",
    )
    if any(kw in desc_lower for kw in code_keywords):
        s = _gather_code_inventory(project_path)
        if s:
            sections.append(s)

    # ─── 6. Config summary ───
    config_keywords = (
        "설정", "config", "configuration", "환경 변수",
        "dockerfile", "docker", "ci", "github actions",
        "build config", "빌드 설정", "deploy", "배포",
    )
    if any(kw in desc_lower for kw in config_keywords):
        s = _gather_config_files(project_path)
        if s:
            sections.append(s)

    if not sections:
        return ""

    return "\n\n".join(sections)


# ─────────── Source-of-truth data gatherers ───────────

# 빌드/캐시 폴더 — 트리 스캔 시 항상 제외
_IGNORE_DIRS = frozenset({
    ".git", "node_modules", "__pycache__", ".venv", "venv",
    ".pytest_cache", ".hypothesis", "dist", "build",
    ".rag_cache", "coverage", ".generated", ".cache",
    ".next", ".turbo", "target", "vendor", "out",
})
_IGNORE_FILES = frozenset({".DS_Store", ".pyc"})


def _gather_folder_tree(project_path: str) -> str:
    """실제 폴더 트리 + 통계 (최대 4 depth, 디렉토리당 30 항목)."""
    try:
        lines = []
        max_depth = 4
        max_entries_per_dir = 30
        root_name = os.path.basename(project_path.rstrip("/")) or project_path

        total_dirs = 0
        total_files = 0
        ext_count = {}

        def _walk(path: str, prefix: str, depth: int):
            nonlocal total_dirs, total_files
            if depth > max_depth:
                return
            try:
                entries = sorted(os.listdir(path))
            except (PermissionError, OSError):
                return
            entries = [e for e in entries if e not in _IGNORE_FILES and e not in _IGNORE_DIRS]
            entries = entries[:max_entries_per_dir]
            for i, e in enumerate(entries):
                full = os.path.join(path, e)
                is_last = (i == len(entries) - 1)
                connector = "└── " if is_last else "├── "
                if os.path.isdir(full):
                    total_dirs += 1
                    lines.append(f"{prefix}{connector}{e}/")
                    next_prefix = prefix + ("    " if is_last else "│   ")
                    _walk(full, next_prefix, depth + 1)
                else:
                    total_files += 1
                    ext = os.path.splitext(e)[1].lstrip(".").lower()
                    if ext:
                        ext_count[ext] = ext_count.get(ext, 0) + 1
                    lines.append(f"{prefix}{connector}{e}")

        lines.append(f"{root_name}/")
        _walk(project_path, "", 1)

        tree_text = "\n".join(lines[:200])
        ext_table_rows = sorted(ext_count.items(), key=lambda x: -x[1])[:10]
        ext_table = "\n".join(f"| .{ext} | {cnt} |" for ext, cnt in ext_table_rows)

        return f"""## 실제 프로젝트 폴더 구조 (디스크 검증)

프로젝트 루트: `{project_path}`

### 디렉토리 트리

```
{tree_text}
```

### 통계

| 항목 | 값 |
|------|-----|
| 총 디렉토리 | {total_dirs} |
| 총 파일 | {total_files} |
| 최대 깊이 | {max_depth} |

### 파일 형식별 분포

| 확장자 | 파일 수 |
|--------|---------|
{ext_table}"""
    except Exception as e:
        print(f"[RealContext/folder] 실패: {e}")
        return ""


def _gather_dependencies(project_path: str) -> str:
    """의존성 매니페스트 파일들에서 패키지 + 버전 추출."""
    try:
        out_blocks = []

        # package.json (npm/yarn)
        pkg = os.path.join(project_path, "package.json")
        if os.path.isfile(pkg):
            try:
                with open(pkg, "r", encoding="utf-8") as f:
                    data = json.load(f)
                deps = data.get("dependencies", {})
                dev_deps = data.get("devDependencies", {})
                rows = []
                for name, ver in sorted(deps.items()):
                    rows.append(f"| {name} | {ver} | runtime |")
                for name, ver in sorted(dev_deps.items()):
                    rows.append(f"| {name} | {ver} | dev |")
                if rows:
                    out_blocks.append(
                        "### Node.js (package.json)\n\n"
                        f"프로젝트명: `{data.get('name', 'N/A')}` · 버전: `{data.get('version', 'N/A')}`\n\n"
                        "| 패키지 | 버전 | 종류 |\n"
                        "|--------|------|------|\n"
                        + "\n".join(rows[:50])
                    )
            except (json.JSONDecodeError, OSError) as e:
                print(f"[RealContext/deps] package.json 파싱 실패: {e}")

        # requirements.txt (pip)
        req = os.path.join(project_path, "requirements.txt")
        if os.path.isfile(req):
            try:
                with open(req, "r", encoding="utf-8") as f:
                    lines = [ln.strip() for ln in f if ln.strip() and not ln.strip().startswith("#")]
                if lines:
                    rows = "\n".join(f"| {ln} |" for ln in lines[:50])
                    out_blocks.append(
                        "### Python (requirements.txt)\n\n"
                        "| 패키지 |\n"
                        "|--------|\n"
                        f"{rows}"
                    )
            except OSError as e:
                print(f"[RealContext/deps] requirements.txt 실패: {e}")

        # ai_engine/requirements.txt (별도 위치)
        req2 = os.path.join(project_path, "ai_engine", "requirements.txt")
        if os.path.isfile(req2) and req2 != req:
            try:
                with open(req2, "r", encoding="utf-8") as f:
                    lines = [ln.strip() for ln in f if ln.strip() and not ln.strip().startswith("#")]
                if lines:
                    rows = "\n".join(f"| {ln} |" for ln in lines[:50])
                    out_blocks.append(
                        "### Python (ai_engine/requirements.txt)\n\n"
                        "| 패키지 |\n"
                        "|--------|\n"
                        f"{rows}"
                    )
            except OSError:
                pass

        # pyproject.toml
        pyproj = os.path.join(project_path, "pyproject.toml")
        if os.path.isfile(pyproj):
            try:
                with open(pyproj, "r", encoding="utf-8") as f:
                    content = f.read()[:3000]
                out_blocks.append(f"### Python (pyproject.toml)\n\n```toml\n{content}\n```")
            except OSError:
                pass

        # Cargo.toml (Rust)
        cargo = os.path.join(project_path, "Cargo.toml")
        if os.path.isfile(cargo):
            try:
                with open(cargo, "r", encoding="utf-8") as f:
                    content = f.read()[:3000]
                out_blocks.append(f"### Rust (Cargo.toml)\n\n```toml\n{content}\n```")
            except OSError:
                pass

        # go.mod
        gomod = os.path.join(project_path, "go.mod")
        if os.path.isfile(gomod):
            try:
                with open(gomod, "r", encoding="utf-8") as f:
                    content = f.read()[:2000]
                out_blocks.append(f"### Go (go.mod)\n\n```\n{content}\n```")
            except OSError:
                pass

        if not out_blocks:
            return ""
        return "## 실제 프로젝트 의존성 (manifest 파일에서 직접 추출)\n\n" + "\n\n".join(out_blocks)
    except Exception as e:
        print(f"[RealContext/deps] 실패: {e}")
        return ""


def _gather_git_summary(project_path: str) -> str:
    """git log + git status 요약."""
    try:
        if not os.path.isdir(os.path.join(project_path, ".git")):
            return ""
        results = []

        # 최근 20개 커밋
        try:
            log_out = subprocess.check_output(
                ["git", "-C", project_path, "log", "--oneline", "-20", "--no-decorate"],
                stderr=subprocess.DEVNULL, timeout=5,
            ).decode("utf-8", errors="ignore").strip()
            if log_out:
                results.append(f"### 최근 커밋 20개\n\n```\n{log_out}\n```")
        except (subprocess.SubprocessError, OSError) as e:
            print(f"[RealContext/git] log 실패: {e}")

        # 현재 브랜치
        try:
            branch = subprocess.check_output(
                ["git", "-C", project_path, "branch", "--show-current"],
                stderr=subprocess.DEVNULL, timeout=5,
            ).decode("utf-8", errors="ignore").strip()
            if branch:
                results.append(f"### 현재 브랜치\n\n`{branch}`")
        except (subprocess.SubprocessError, OSError):
            pass

        # git status (수정/staged 파일)
        try:
            status = subprocess.check_output(
                ["git", "-C", project_path, "status", "--short"],
                stderr=subprocess.DEVNULL, timeout=5,
            ).decode("utf-8", errors="ignore").strip()
            if status:
                results.append(f"### 미커밋 변경사항\n\n```\n{status[:1500]}\n```")
        except (subprocess.SubprocessError, OSError):
            pass

        # 최근 변경 파일 빈도 (1주일)
        try:
            files = subprocess.check_output(
                ["git", "-C", project_path, "log", "--name-only", "--pretty=format:",
                 "--since=1.week"],
                stderr=subprocess.DEVNULL, timeout=10,
            ).decode("utf-8", errors="ignore").strip()
            if files:
                from collections import Counter
                counter = Counter(line for line in files.splitlines() if line.strip())
                top = counter.most_common(15)
                rows = "\n".join(f"| {f} | {c} |" for f, c in top)
                results.append(
                    "### 최근 1주일 변경 빈도 상위\n\n"
                    "| 파일 | 변경 횟수 |\n"
                    "|------|---------|\n"
                    f"{rows}"
                )
        except (subprocess.SubprocessError, OSError):
            pass

        if not results:
            return ""
        return "## 실제 Git 데이터 (저장소에서 직접 추출)\n\n" + "\n\n".join(results)
    except Exception as e:
        print(f"[RealContext/git] 실패: {e}")
        return ""


def _gather_readme(project_path: str) -> str:
    """README/CHANGELOG 본문 추출."""
    try:
        out_blocks = []
        candidates = ["README.md", "README.rst", "README.txt", "README",
                      "CHANGELOG.md", "CHANGELOG.txt", "CHANGES.md"]
        for name in candidates:
            full = os.path.join(project_path, name)
            if not os.path.isfile(full):
                # case-insensitive 검색 (한 번만)
                try:
                    actual = next(
                        (e for e in os.listdir(project_path)
                         if e.lower() == name.lower()),
                        None,
                    )
                    if actual:
                        full = os.path.join(project_path, actual)
                    else:
                        continue
                except OSError:
                    continue
            if os.path.isfile(full):
                try:
                    with open(full, "r", encoding="utf-8", errors="ignore") as f:
                        content = f.read()
                    # 길이 제한
                    if len(content) > 4000:
                        content = content[:4000] + "\n\n... (생략)"
                    out_blocks.append(f"### {os.path.basename(full)}\n\n{content}")
                except OSError as e:
                    print(f"[RealContext/readme] {name} 실패: {e}")
        if not out_blocks:
            return ""
        return "## 실제 README/CHANGELOG (디스크에서 직접 읽음)\n\n" + "\n\n".join(out_blocks)
    except Exception as e:
        print(f"[RealContext/readme] 실패: {e}")
        return ""


def _gather_code_inventory(project_path: str) -> str:
    """언어별 파일 개수 + LOC 상위 10개 파일."""
    try:
        # 언어별 카운트
        lang_count = {}
        lang_loc = {}
        # LOC 상위
        file_loc = []  # [(loc, relpath)]

        ext_to_lang = {
            "py": "Python", "js": "JavaScript", "jsx": "JavaScript",
            "ts": "TypeScript", "tsx": "TypeScript", "java": "Java",
            "go": "Go", "rs": "Rust", "rb": "Ruby", "php": "PHP",
            "c": "C", "h": "C", "cpp": "C++", "hpp": "C++",
            "cs": "C#", "swift": "Swift", "kt": "Kotlin",
            "html": "HTML", "css": "CSS", "scss": "SCSS",
            "md": "Markdown", "yaml": "YAML", "yml": "YAML",
            "json": "JSON", "sh": "Shell", "sql": "SQL",
        }

        for root, dirs, files in os.walk(project_path):
            # 빌드/캐시 폴더 제외
            dirs[:] = [d for d in dirs if d not in _IGNORE_DIRS]
            for f in files:
                if f in _IGNORE_FILES:
                    continue
                ext = os.path.splitext(f)[1].lstrip(".").lower()
                lang = ext_to_lang.get(ext)
                if not lang:
                    continue
                full = os.path.join(root, f)
                try:
                    with open(full, "r", encoding="utf-8", errors="ignore") as fh:
                        loc = sum(1 for _ in fh)
                except OSError:
                    continue
                lang_count[lang] = lang_count.get(lang, 0) + 1
                lang_loc[lang] = lang_loc.get(lang, 0) + loc
                rel = os.path.relpath(full, project_path)
                file_loc.append((loc, rel))

        if not lang_count:
            return ""

        lang_rows = "\n".join(
            f"| {lang} | {lang_count[lang]} | {lang_loc.get(lang, 0):,} |"
            for lang in sorted(lang_count.keys(), key=lambda l: -lang_count[l])[:15]
        )

        file_loc.sort(reverse=True)
        top_rows = "\n".join(f"| {rel} | {loc:,} |" for loc, rel in file_loc[:10])

        return f"""## 실제 코드 인벤토리 (디스크 스캔)

### 언어별 통계

| 언어 | 파일 수 | 총 라인 |
|------|--------|--------|
{lang_rows}

### 라인 수 상위 10개 파일

| 파일 | LOC |
|------|-----|
{top_rows}"""
    except Exception as e:
        print(f"[RealContext/inventory] 실패: {e}")
        return ""


def _gather_config_files(project_path: str) -> str:
    """주요 설정 파일 목록 + 일부 내용."""
    try:
        out_blocks = []
        config_files = [
            ".gitignore", "Dockerfile", "docker-compose.yml", "docker-compose.yaml",
            "Makefile", ".env.example", "tsconfig.json", "vite.config.js",
            "webpack.config.js", "next.config.js", "tailwind.config.js",
        ]
        gh_actions = os.path.join(project_path, ".github", "workflows")

        found = []
        for name in config_files:
            full = os.path.join(project_path, name)
            if os.path.isfile(full):
                try:
                    size = os.path.getsize(full)
                    found.append((name, full, size))
                except OSError:
                    pass

        # GitHub Actions
        if os.path.isdir(gh_actions):
            try:
                for f in os.listdir(gh_actions):
                    full = os.path.join(gh_actions, f)
                    if os.path.isfile(full):
                        size = os.path.getsize(full)
                        found.append((f".github/workflows/{f}", full, size))
            except OSError:
                pass

        if not found:
            return ""

        rows = "\n".join(f"| `{name}` | {size} bytes |" for name, _, size in found[:20])
        out_blocks.append(
            "### 발견된 설정 파일\n\n"
            "| 파일 | 크기 |\n"
            "|------|-----|\n"
            f"{rows}"
        )

        # 핵심 파일은 내용 일부 첨부
        for name, full, size in found[:3]:
            if size > 4000:
                continue  # 너무 큰 파일은 스킵
            try:
                with open(full, "r", encoding="utf-8", errors="ignore") as f:
                    content = f.read()[:1500]
                out_blocks.append(f"### `{name}` 내용\n\n```\n{content}\n```")
            except OSError:
                pass

        return "## 실제 설정 파일 (디스크에서 직접 읽음)\n\n" + "\n\n".join(out_blocks)
    except Exception as e:
        print(f"[RealContext/config] 실패: {e}")
        return ""


def _gather_real_context_forced(project_path: str, title: str = "") -> str:
    """키워드 매칭 없이 forced fallback 시 강제로 실제 데이터 수집.

    description에 명시된 키워드가 없어도 폴더 트리 + 코드 인벤토리는 항상 유용한
    데이터이므로 fallback 경로에서는 무조건 활성화.
    """
    if not project_path or not os.path.isdir(project_path):
        return ""
    sections = []
    folder = _gather_folder_tree(project_path)
    if folder:
        sections.append(folder)
    inventory = _gather_code_inventory(project_path)
    if inventory:
        sections.append(inventory)
    if not sections:
        return ""
    return "\n\n".join(sections)


async def _enrich_content_via_gateway(
    gw, model_id: str, primary_tool: str, title: str,
    description: str, final_text: str, max_tokens: int = 2000,
    project_path: str = "",
) -> str:
    """게이트웨이 호출로 파일 본문 콘텐츠 보강 — 퀄리티 보장.

    forced fallback 직전에 호출해서 빈약한 final_text를 풍부한 본문으로 확장.
    형식별 최소 분량 임계치를 만족할 때까지 최대 2회 재시도.
    게이트웨이를 반드시 경유 (비용/사용량 측정 유지).

    Args:
        gw: GatewayClient instance
        model_id: 호출할 모델 ID (이미 자동 라우팅된 Claude 권장)
        primary_tool: 작업 종류 (generate_pdf 등)
        title, description, final_text: 원본 컨텍스트

    Returns:
        보강된 마크다운 텍스트. 모든 시도 실패 시 final_text 그대로 반환.
    """
    pt = (primary_tool or "").lower()
    # 형식별 퀄리티 임계치 — 사용자가 보고서로 사용해도 충분한 분량
    quality_min = {
        "generate_pdf":   1500,
        "generate_docx":  1500,
        "generate_pptx":   800,   # 슬라이드는 압축적
        "generate_xlsx":   400,   # 표는 데이터 위주
        "generate_image":  150,   # 이미지 prompt
        "write_file":      400,
    }.get(pt, 600)

    style_hint = {
        "generate_pdf":  "다단 마크다운 보고서 형식 (## 헤더, 본문 단락, 필요 시 표).",
        "generate_pptx": "슬라이드 헤더와 불릿 위주 (## 섹션 = 슬라이드 제목, 각 섹션 4-6개 불릿).",
        "generate_xlsx": "마크다운 표 형식 우선 (| 헤더 | ... |), 표 없으면 카테고리별 데이터.",
        "generate_docx": "다단 마크다운 문서 (## 섹션, 본문 단락, 불릿 리스트).",
        "generate_image": "이미지 생성용 영문 prompt 한 줄 (구체적 스타일/구성/조명/시점).",
        "write_file":    "코드/마크다운/텍스트 콘텐츠.",
    }.get(pt, "마크다운 형식의 충실한 본문.")

    # 형식별 추가 퀄리티 가이드라인
    quality_guide = {
        "generate_pdf":   "최소 4개 섹션, 각 섹션 2-4 문단의 깊이 있는 설명. Lorem ipsum이나 placeholder 금지 — 실제 정보로 채울 것.",
        "generate_docx":  "최소 4개 섹션, 각 섹션 2-4 문단의 풍부한 내용. 도입→본론→결론 구조. 실제 정보 위주.",
        "generate_pptx":  "최소 5개 슬라이드(권장 7-10개). 각 슬라이드 3-5개 구체적 불릿. 첫 슬라이드 개요, 마지막 결론/Next Steps.",
        "generate_xlsx":  "헤더 행 + 최소 5행 이상의 의미있는 데이터. 가능하면 10-20행. 합리적인 실제 값.",
        "generate_image": "스타일/구성/조명/색감/시점을 모두 포함한 구체적 영문 prompt 1줄.",
        "write_file":     "실제 사용 가능한 풍부한 콘텐츠. 빈 placeholder 금지.",
    }.get(pt, "충실하고 구체적인 콘텐츠로 작성.")

    sys_prompt = f"""당신은 파일 생성 작업의 콘텐츠 작성자입니다.
주제와 지시사항에 맞는 충실한 본문을 작성하세요.

[출력 형식] {style_hint}
[퀄리티 기준] {quality_guide}
[제약] 도구 호출은 하지 않습니다. 마크다운 본문만 출력하세요.
[중요] '생성 완료' 같은 거짓 주장 금지. 실제 콘텐츠만 작성.
"""

    # === 약점 2 개선: 실제 디스크 데이터 주입 ===
    # 폴더 구조/파일 목록 같은 사실 기반 작업은 LLM이 추측하지 않도록 실제 데이터 첨부.
    real_context = _gather_real_context(description, project_path)

    user_msg_base = f"""작업: {title}
지시사항: {description}

{f'### 실제 프로젝트 데이터 (디스크에서 직접 수집)\n\n{real_context}\n\n위 실제 데이터를 그대로 인용해서 작성하세요. 추측이나 가공 금지.\n' if real_context else ''}
이전 응답 일부 (참고):
{(final_text or '')[:1500]}

위 정보를 바탕으로 [{pt}] 작업에 사용할 본문을 마크다운으로 작성해주세요. 최소 {quality_min}자 이상의 풍부한 내용을 작성하세요."""

    # 최대 2회 시도 — 첫 시도 결과가 임계치 미만이면 재시도
    best = ""
    for attempt in range(2):
        user_msg = user_msg_base
        if attempt > 0 and best:
            # 두 번째 시도 — 첫 결과보다 더 풍부하게 작성하라고 지시
            user_msg += f"\n\n[재시도] 이전 시도는 분량이 부족했습니다 ({len(best)}자). 더 깊이 있고 풍부하게 작성해주세요."
        messages = [{"role": "user", "content": [{"text": user_msg}]}]
        try:
            result = await asyncio.wait_for(
                gw.converse(model_id=model_id, messages=messages, system_prompt=sys_prompt),
                timeout=90,
            )
            if result.get("decision") != "ALLOW":
                print(f"[Enrich] gateway 거부(attempt {attempt+1}): {result.get('error') or result.get('decision')}")
                continue
            output = result.get("output", {}).get("message", {}).get("content", [])
            text = "\n".join(c.get("text", "") for c in output if "text" in c).strip()
            if text and len(text) > len(best):
                best = text
            if best and len(best) >= quality_min:
                # 퀄리티 임계치 만족 — 즉시 반환
                if attempt > 0:
                    print(f"[Enrich] 재시도로 퀄리티 임계치 달성 ({len(best)}자 >= {quality_min}자)")
                return best
        except Exception as e:
            print(f"[Enrich] 예외(attempt {attempt+1}): {e}")
            continue

    # 임계치 미달이지만 가지고 있는 best가 final_text보다 풍부하면 그것 사용
    if best and len(best) > len(final_text or ""):
        print(f"[Enrich] 임계치({quality_min}) 미달이지만 best({len(best)}자) 사용")
        return best
    return final_text or description or title or "내용 없음"


def _prune_orphan_generated_images(project_path: str, max_age_seconds: int = 3600) -> int:
    """Remove standalone PNG files in .generated/ older than max_age_seconds.

    "Orphan" definition: a .png whose mtime is older than the cutoff AND
    has no matching parent document (.pdf/.pptx/.docx/.xlsx) created within
    a 60-second window of it.

    Returns the number of files deleted. Conservative — preserves anything
    that might still be referenced by an ongoing or recent generation run.
    """
    import time as _t
    try:
        local_root = _resolve_local_root(project_path)
        gen_dir = os.path.join(local_root, ".generated")
        if not os.path.isdir(gen_dir):
            return 0
    except Exception:
        return 0

    now = _t.time()
    cutoff = now - max_age_seconds

    # Index .pdf/.pptx/.docx/.xlsx mtimes for pairing check.
    doc_mtimes = []
    for entry in os.scandir(gen_dir):
        if not entry.is_file():
            continue
        ext = entry.name.rsplit(".", 1)[-1].lower()
        if ext in ("pdf", "pptx", "docx", "xlsx"):
            try:
                doc_mtimes.append(entry.stat().st_mtime)
            except OSError:
                continue

    deleted = 0
    for entry in os.scandir(gen_dir):
        if not entry.is_file():
            continue
        name = entry.name
        if not name.lower().endswith((".png", ".jpg", ".jpeg", ".webp")):
            continue
        try:
            st = entry.stat()
        except OSError:
            continue
        if st.st_mtime > cutoff:
            continue  # too recent — keep
        # Only delete if no doc was created within ±60s of this image.
        nearby_doc = any(abs(dm - st.st_mtime) <= 60 for dm in doc_mtimes)
        if nearby_doc:
            continue
        try:
            os.remove(entry.path)
            # Also remove sidecar meta if present
            meta_path = entry.path + ".meta.json"
            if os.path.isfile(meta_path):
                try:
                    os.remove(meta_path)
                except OSError:
                    pass
            deleted += 1
        except OSError as e:
            print(f"[ForceGenerate] orphan prune skip {name}: {e}")
    if deleted:
        print(f"[ForceGenerate] orphan prune: {deleted} stale image(s) removed from {gen_dir}")
    return deleted


def _resolve_active_template(template_id, store_root):
    """활성 템플릿을 1회 해석한다 (요구사항 5.1, 5.3, 5.4, 5.5). 설계 §구성요소 6.

    Returns:
        (template_path, style_profile, used) 튜플.
        - template_id 없음/"" → (None, None, False)  # 요구사항 5.2 (무템플릿 — 기존 동작 보존)
        - template-not-found 또는 get_template 에러 → 로그(≤200자) + (None, None, False)  # 요구사항 5.4
        - base.pptx 부재/비파일 또는 styleProfile 부재/로드 실패 → 로그(≤200자) + (None, None, False)  # 요구사항 5.5
        - 정상 → (abs_templatePath, styleProfile_dict, True)  # 요구사항 5.3
    어떤 경우에도 예외를 밖으로 던지지 않는다(폴백 격리, 요구사항 9). 실패는 항상 무템플릿 경로로 귀결된다.
    """
    # 요구사항 5.2 — templateId 미전달/"" 이면 무템플릿. get_template 조차 호출하지 않는다.
    if not template_id:
        return (None, None, False)
    try:
        tm = _load_template_manager()
        if tm is None:
            # template_manager 모듈 부재 — 무템플릿 진행 (요구사항 5.4/5.5 정신과 동일, 격리)
            print(f"[ActiveTemplate] template_manager 사용 불가 → 무템플릿 진행 (templateId={str(template_id)[:80]})")
            return (None, None, False)

        result = tm.get_template(template_id, store_root)
        if not isinstance(result, dict) or result.get("error"):
            # template-not-found / invalid-template-id 등 — 무템플릿 경로로 진행 (요구사항 5.4)
            cause = (result.get("error") if isinstance(result, dict) else str(result)) or "unknown"
            print(f"[ActiveTemplate] get_template 실패 → no-template fallback: {str(cause)[:200]}")
            return (None, None, False)

        # 기준 .pptx 경로 검증 — 부재/비파일이면 무템플릿 진행 (요구사항 5.5)
        template_path = result.get("templatePath") or ""
        if not template_path or not os.path.isfile(template_path):
            print(f"[ActiveTemplate] base.pptx 로드 실패 → no-template fallback: {str(template_path)[:200]}")
            return (None, None, False)

        # Style_Profile 검증 — dict 가 아니거나 error 필드가 있으면 로드 실패로 간주 (요구사항 5.5)
        style_profile = result.get("styleProfile")
        if not isinstance(style_profile, dict) or style_profile.get("error"):
            err = style_profile.get("error") if isinstance(style_profile, dict) else "missing"
            print(f"[ActiveTemplate] style_profile 로드 실패 → no-template fallback: {str(err)[:200]}")
            return (None, None, False)

        # 정상 — 절대 경로 + Style_Profile dict 반환 (요구사항 5.3)
        return (os.path.abspath(template_path), style_profile, True)
    except Exception as e:
        # 어떤 예외도 생성 경로로 전파하지 않는다 — 무템플릿으로 안전 폴백 (요구사항 5.5, 9)
        print(f"[ActiveTemplate] 예외 → no-template fallback: {str(e)[:200]}")
        return (None, None, False)


def _infer_file_intent_from_prompt(prompt: str, client_intent: str = "", final_text: str = ""):
    """단일 호출(run-agent)에서 사용자가 파일 생성을 원했는지 추론한다 (이슈 2).

    오케스트레이터의 강한/약한 키워드 판정 로직과 동일한 원칙을 프롬프트에 직접 적용:
    - 강한 키워드(pdf/pptx/xlsx/docx/png 등 명시적 형식, "보고서/발표자료" 등)가
      있으면 파일 생성 의도로 간주하고 primary_tool을 추론한다.
    - 약한 키워드(분석/리뷰/검토 등)만 있으면 텍스트 응답으로 충분하다고 보고
      강제 생성을 발동하지 않는다 (오탐 방지).
    - client_intent == "file_generation"이면 강제로 wanted=True.

    Returns:
        (primary_tool, wanted_files, target_files)
        primary_tool: "generate_pdf"|"generate_pptx"|... (wanted=False면 "")
        wanted_files: bool
        target_files: list (보통 빈 리스트 — _force_generate_from_text가 ext로 결정)
    """
    text_lower = (prompt or "").lower()
    primary_tool = ""
    # 이미지 생성 의도 — 명시적 형식 문서보다 먼저 체크 (이슈 1).
    # 편집/분석/수정이 아닌 순수 "생성/그리기" 의도일 때만 generate_image.
    _img_gen = bool(re.search(
        r"이미지\s*(생성|만들|그려|그리)|이미지로\s*(만들|생성|그려|변환)|"
        r"사진.*만들|일러스트|로고.*(디자인|만들)|배너.*만들|아이콘.*만들|썸네일|"
        r"(create|generate)\s+(an?\s+)?(image|picture|illustration|logo|banner|icon)",
        text_lower,
    ))
    _img_edit = bool(re.search(
        r"(분석|설명|읽어|읽기|수정|편집|지워|제거|배경|inpaint|outpaint)", text_lower
    ))
    if _img_gen and not _img_edit:
        return "generate_image", True, []

    # 명시적 형식/산출물 키워드 — 이게 있어야 강제 생성 진입
    strong_map = [
        (("pptx", "파워포인트", "프레젠테이션", "발표자료", "발표 자료", "슬라이드", "deck", "ppt"), "generate_pptx"),
        (("xlsx", "엑셀", "스프레드시트", "spreadsheet", "excel"), "generate_xlsx"),
        (("docx", "워드", "word 문서", "워드문서"), "generate_docx"),
        (("pdf", "보고서", "리포트", "report"), "generate_pdf"),
    ]
    has_strong = False
    for kws, tool in strong_map:
        if any(kw in text_lower for kw in kws):
            primary_tool = tool
            has_strong = True
            break

    # client_intent 힌트 — file_generation이면 형식 미지정이어도 pdf로 폴백
    intent_wants = (client_intent or "").lower() == "file_generation"

    wanted = has_strong or intent_wants
    if wanted and not primary_tool:
        primary_tool = "generate_pdf"
    if not wanted:
        return "", False, []
    return primary_tool, True, []


async def _force_generate_from_text(
    primary_tool: str,
    target_files: list,
    title: str,
    description: str,
    final_text: str,
    project_path: str,
    aws_profile: str,
    bedrock_user: str,
    template_id: str = "",
):
    """결정적 도구 디스패처 — final_text(보통 _enrich_content_via_gateway로 보강됨)를
    파일로 변환한다.

    퀄리티 보장:
    - 호출 전에 _enrich_content_via_gateway로 Claude가 풍부한 콘텐츠를 만들어 둠
    - 이 함수는 마크다운 → 결정적 라이브러리(reportlab/python-pptx/openpyxl/python-docx) 변환만 수행
    - 따라서 결과 파일의 내용 퀄리티 = Claude 출력 퀄리티 그대로

    primary_tool과 target_files의 확장자를 보고 어떤 도구를 호출할지 결정한다.

    Returns:
        list of (relative_path, file_info_dict) tuples for verified files.
    """
    pt = (primary_tool or "").lower()
    # 활성 템플릿 1회 해석 (요구사항 5.1, 5.3, 5.4, 5.5) — 설계 §구성요소 6.
    # templateId가 없거나 무효하면 (_tpl_path, _tpl_profile, _tpl_used) = (None, None, False)로
    # 무템플릿 경로(기존 동작)가 바이트 단위로 보존된다(요구사항 5.2). 유효하면 아래 PPTX 분기에서
    # inp에 templatePath/templateId/styleProfile을 주입한다. _tpl_profile은 11.x(배경 Tier 토큰
    # 주입)에서 동일 함수 내 로컬로 재사용된다. store_root는 1회만 결정한다(1회 로드).
    _tpl_store_root = None
    try:
        _tm_for_store = _load_template_manager()
        if _tm_for_store is not None:
            _tpl_store_root = _tm_for_store.resolve_template_store_root()
    except Exception as _e:
        # store_root 결정 실패는 무템플릿으로 안전 폴백 (요구사항 5.5, 9)
        print(f"[ForceGenerate] template store_root 결정 실패 → 무템플릿: {str(_e)[:200]}")
        _tpl_store_root = None
    _tpl_path, _tpl_profile, _tpl_used = _resolve_active_template(template_id, _tpl_store_root)
    # matplotlib(Tier 2) 팔레트 1회 계산 (요구사항 7.3). _tpl_profile이 dict이고 primary/
    # secondary/accent 중 유효 #RRGGBB가 2색 이상이면 [primary, ...] 리스트, 아니면 None.
    # None이면 _tool_generate_native_diagram이 기존 기본 색상으로 폴백한다(요구사항 7.5/5.2).
    _tpl_palette = _build_palette(_tpl_profile)
    # target_files에서 확장자 결정 → 없으면 primary_tool로 결정
    needed_exts = []
    if target_files:
        for tf in target_files:
            ext = (str(tf).split(".")[-1] or "").lower()
            if ext and ext not in needed_exts:
                needed_exts.append(ext)
    if not needed_exts:
        # primary_tool로 매핑
        tool_to_ext = {
            "generate_pdf": "pdf",
            "generate_pptx": "pptx",
            "generate_xlsx": "xlsx",
            "generate_docx": "docx",
            "generate_image": "png",
            "write_file": "md",
            # 코드 분석/리팩토링도 결과물을 마크다운 보고서로 저장 — 빈손 종료 방지
            "code": "md",
            "run_command": "md",
        }
        if pt in tool_to_ext:
            needed_exts.append(tool_to_ext[pt])
        else:
            # 알 수 없는 primary_tool — 일단 마크다운으로라도 저장
            needed_exts.append("md")

    if not needed_exts:
        return []

    # 본문 텍스트 — final_text가 비어있으면 description, title 순으로 fallback
    base_text = (final_text or "").strip() or (description or "").strip() or (title or "문서")

    # 섹션 단위로 분리 (마크다운 헤더 + 빈 줄)
    sections = _split_into_sections(base_text, fallback_title=title or "본문")

    # 시각 콘텐츠 요구 감지 — description/title/final_text 어느 곳에서든 키워드 등장 시 True
    # PDF/PPTX/DOCX의 각 섹션에 자동으로 imagePrompt를 부여하여 실제 이미지 임베드를 유도한다.
    visual_intent = (
        _detect_visual_intent(description)
        or _detect_visual_intent(title)
        or _detect_visual_intent(final_text or "")
    )

    # === 이미지 생성 우선순위 ===
    # tier 0:    HTML 슬라이드 (Genspark/Gamma 수준) — Electron bridge 활성 시
    # tier 0.6:  Vertex AI 이미지 (Nano Banana Pro / Imagen 4, AE_ENABLE_VERTEX_IMAGE=1)
    #            — 텍스트 정확 렌더, Bedrock 동급 모델 부재 시 사용자 결정으로 예외 허용.
    # tier 0.7:  Hero photo (Bedrock Stability Ultra, AE_BEDROCK_HERO_IMAGE=1, 표지 한정)
    # tier 1:    Mermaid 다이어그램 (Claude generates code via Gateway, Korean-accurate)
    # tier 2:    matplotlib native diagram (디렉토리/플로우/표 — 결정론적)
    # tier 3:    python-pptx 네이티브 텍스트 슬라이드 (이미지 없는 경우)
    #
    # Tier 0.5 (Bedrock generative image for slide bodies) is OFF by default —
    # generative models cannot render reliable text, producing "fake glyph" output
    # (Korean: 킹딕록켠; English: TEUQUB, PEDUCTKIN) that misaligns with slide
    # content. Re-enable only for debugging via AE_ENABLE_BEDROCK_SLIDE_IMAGES=1.
    #
    # LLM 호출(Mermaid 코드 생성, operation JSON 등)은 Bedrock Gateway 경유
    # (gateway.md 규칙). 이미지 생성에 한해 Bedrock에 동급 모델이 없는 현실을
    # 반영해 Vertex AI 호출을 사용자 결정으로 예외 허용 (Tier 0.6). Bedrock
    # 이미지 모델(Tier 0.5/0.7)도 Gateway 경유 그대로 유지. 직접 boto3/Anthropic
    # SDK 사용 금지.
    #
    # AE_FORCE_NATIVE_DIAGRAM=1 환경변수로 디버깅 시 native만 강제 가능.
    force_native = os.environ.get("AE_FORCE_NATIVE_DIAGRAM", "").strip() == "1"

    # Mermaid + matplotlib 경로는 visual_intent가 있는 모든 섹션에 대해 활성화.
    # 이전: structural 키워드(폴더/트리/플로우)가 있을 때만 진입 → 비구조 섹션은
    #   slide body가 비어 있고 Bedrock이 fake-glyph를 채워 넣었음.
    # 이제: visual_intent면 항상 mermaid 시도 → 실패 시 matplotlib classifier
    #   fallback. HTML 배경이 있는 섹션은 mermaid 태스크 필터에서 자동 skip.
    # 편집 가능 우선 모드(AE_PREFER_EDITABLE_DIAGRAM!=0, 기본)에서는 mermaid/matplotlib
    # '통짜 PNG' 다이어그램 생성을 끈다 — 대신 _tool_generate_pptx가 LLM 구조화로
    # 편집 가능한 네이티브 도형을 그린다(사용자 요구: 모든 도형 개별 편집 가능).
    _editable_pref = os.environ.get("AE_PREFER_EDITABLE_DIAGRAM", "1") != "0"
    use_native = bool(visual_intent) and not _editable_pref
    # force_native / circuit-broken 은 진단용으로만 의미가 있고, 실제 진입 조건
    # 자체는 visual_intent로 단순화. (force_native 변수는 하위 로깅에서 사용 가능.)
    _ = force_native  # noqa: F841 — kept for future diagnostic logging

    # Tier ordering smoke log — surfaces routing decision in production logs so
    # we can confirm Bedrock image gen is OFF by default (no fake-glyph output).
    _bedrock_slide_on = os.environ.get("AE_ENABLE_BEDROCK_SLIDE_IMAGES", "").strip() == "1"
    _hero_on = os.environ.get("AE_BEDROCK_HERO_IMAGE", "").strip() == "1"
    # Vertex AI tier 0.6 — auto-enabled when a key is resolvable from
    # AWS Secrets Manager / env / local cache. Probe once per call so the
    # tier-ordering log reflects current availability for this user.
    try:
        from ai_engine.vertex_image_module import get_vertex_image_client as _vget
        _vertex_probe = _vget(aws_profile=aws_profile or "")
        _vertex_on = bool(_vertex_probe.enabled)
    except Exception:
        _vertex_on = False
    print(
        "[ForceGenerate] tier ordering: HTML(if bridge), "
        f"{'Vertex(NanoBanana) ON, ' if _vertex_on else ''}"
        f"{'Hero(Bedrock) ON, ' if _hero_on else ''}"
        "Mermaid, matplotlib, python-pptx native "
        f"(Bedrock slide images {'ON' if _bedrock_slide_on else 'OFF'})"
    )

    # Prune orphan intermediate PNGs that accumulated from previous runs.
    # Keeps recent files (< 1 hour) and any image with a sibling document
    # created within ±60s. Configurable via AE_ORPHAN_TTL_SEC.
    try:
        _orphan_ttl = int(os.environ.get("AE_ORPHAN_TTL_SEC", "3600"))
        _prune_orphan_generated_images(project_path, max_age_seconds=_orphan_ttl)
    except Exception as _e_prune:
        print(f"[ForceGenerate] orphan prune outer error: {_e_prune}")

    # === 섹션별 의미 인식 다이어그램 ===
    # 이전 버그: 모든 섹션에 동일한 shared_native_image를 첨부 → PDF/PPTX 모든 페이지에
    # 똑같은 그림이 반복되고, 그 그림은 마크다운 텍스트를 행 박스로 캡쳐한 거라 도식이 아니었음.
    # 이제 각 섹션의 heading/body를 분석해 적합한 다이어그램 종류를 결정 + 각자 별도 이미지 생성.
    #
    # _section_diagrams: {section_index: relative_path} — _img_attached 분기에서 사용.
    section_diagrams: dict = {}
    # _section_backgrounds: {section_index: relative_path} — full-bleed PPTX bg.
    # PPTX에서 1920×1080 HTML→PNG 캡쳐를 슬라이드 전체 배경으로 사용 (Genspark/Gamma 수준).
    # PDF/DOCX는 imageFile과 동일하게 cap 적용 후 임베드.
    section_backgrounds: dict = {}
    # 문제2 — 네이티브(편집 가능) 다이어그램용 spec 맵. {section_index: {type, content}}
    # ★중요: section_diagrams/section_backgrounds와 동일하게 use_native 블록 *밖*에서
    # 무조건 초기화한다. 이전 버그: 이 변수를 `if use_native:` 안에서만 정의해서,
    # 게이트웨이 실패 등으로 native 다이어그램 단계를 못 타면 PPTX 매핑부의
    # section_diagram_specs.get() 참조가 UnboundLocalError를 던져 전체 PPTX가 0건이 됐다
    # (원격+게이트웨이 403 환경에서 "강제 생성 폴백도 실패" 증상의 직접 원인).
    section_diagram_specs: dict = {}
    shared_native_image: str = ""  # 표지/이미지 ext의 fallback용

    # === HTML Slide Backgrounds (tier 0 — Genspark/Gamma class) =============
    # 이슈 4 — 옵션 A: full-bleed HTML→PNG 배경은 슬라이드 전체를 이미지로 덮어
    # PowerPoint 편집이 불가능하다(텍스트/도형 직접 수정 불가). 사용자 결정으로
    # 기본을 "네이티브 편집 가능 슬라이드"로 전환한다.
    #   - 기본(OFF): HTML 풀블리드 배경을 만들지 않는다 → 텍스트는 placeholder,
    #     다이어그램/이미지는 우측 절반 보조 요소로만 삽입 → 전부 편집 가능.
    #   - AE_ENABLE_HTML_SLIDES=1 로 옵트인하면 기존 풀블리드 고품질 경로 복원.
    # 하위호환: 과거 AE_DISABLE_HTML_SLIDES=1 (끄기) 플래그도 계속 존중한다.
    # 하이브리드 — 기본은 편집 가능 네이티브(OFF). 사용자가 '고품질/PDF용/젠스파크급'을
    # 요청한 경우(title/description/final_text 의도)에 한해 HTML 렌더 경로를 켠다.
    # (AE_ENABLE_HTML_SLIDES=1 명시 시 항상 ON, =0이면 의도와 무관하게 강제 OFF.)
    # 기본 고품질(HTML 렌더). 사용자가 '편집 가능/네이티브'를 명시한 경우에만 네이티브.
    # 의도 판정은 사용자 지시(title/description)만 본다 — 생성 본문(final_text)의
    # 우연한 '수정/편집' 단어로 인한 오탐을 막는다.
    _html_opt_in = _resolve_html_slides(
        os.environ.get("AE_ENABLE_HTML_SLIDES", ""),
        title, description,
    )
    _html_legacy_disabled = os.environ.get("AE_DISABLE_HTML_SLIDES", "").strip() == "1"
    html_slides_enabled = (
        _html_opt_in
        and not _html_legacy_disabled
        and any(e in needed_exts for e in ("pptx", "pdf", "docx"))
    )
    # Bridge 가용성 사전 체크 — Electron이 떠있고 IPC가 활성화돼야 함.
    # Bridge unreachable이면 시도 자체를 skip해서 시간/토큰 절약 + 명확한 로그.
    bridge_available = False
    if html_slides_enabled:
        try:
            _bridge_check = _call_bridge("status", {})
            bridge_available = bool(_bridge_check) or bool(_find_local_chrome())
        except Exception:
            bridge_available = bool(_find_local_chrome())
        if not bridge_available:
            print(
                "[ForceGenerate] HTML slide skip — Electron 브리지/로컬 Chrome 모두 없음. "
                "네이티브 도형 경로로 진행. (Chrome 설치 또는 Electron 앱 실행 시 Genspark급 HTML 렌더)"
            )

    if html_slides_enabled and bridge_available:
        try:
            gw_for_html = _get_gw(aws_profile, bedrock_user)
            html_model = _specialized_model_for_task(
                "file_generation", "",
                aws_profile=aws_profile, bedrock_user=bedrock_user,
            )
        except Exception as e:
            print(f"[ForceGenerate] HTML slide 환경 준비 실패: {e}")
            gw_for_html = None
            html_model = ""

        if gw_for_html and html_model:
            doc_ctx = (title or description or "")[:300]
            print(f"[ForceGenerate] HTML slide tier 0 활성 — model={html_model}, sections={min(len(sections), 8)}")

            async def _try_html_slide_for_section(idx, heading, body):
                # cover slide for first section if visual_intent and a section heading
                # placeholder — keep it simple: every section tries layout pick.
                rel = await _generate_html_slide_for_section(
                    gw_for_html, html_model,
                    section_heading=heading,
                    section_body=body,
                    doc_context=doc_ctx,
                    project_path=project_path,
                )
                return idx, rel

            # Limit to first 8 sections to keep cost/time bounded.
            html_tasks = []
            for idx, sec in enumerate(sections[:8]):
                heading = (sec.get("heading") or "").strip()
                body = (sec.get("body") or "").strip()
                if not (heading or body):
                    continue
                html_tasks.append(_try_html_slide_for_section(idx, heading, body))

            attempted_count = len(html_tasks)
            if html_tasks:
                try:
                    html_results = await asyncio.gather(*html_tasks, return_exceptions=True)
                    fail_count = 0
                    for r in html_results:
                        if isinstance(r, Exception):
                            fail_count += 1
                            print(f"[ForceGenerate] HTML slide task exception: {r}")
                            continue
                        idx, rel = r
                        if rel:
                            section_backgrounds[idx] = rel
                            if not shared_native_image:
                                shared_native_image = rel
                        else:
                            fail_count += 1
                    if fail_count:
                        print(f"[ForceGenerate] HTML slide: {len(section_backgrounds)}/{attempted_count} captured ({fail_count} failed → mermaid/matplotlib fallback로 떨어짐)")
                    else:
                        print(f"[ForceGenerate] HTML slide: {len(section_backgrounds)}/{attempted_count} captured (전체 성공)")
                except Exception as e:
                    print(f"[ForceGenerate] HTML slide gather 실패: {e}")

            if section_backgrounds:
                print(f"[ForceGenerate] HTML slides: {len(section_backgrounds)} captured (Genspark/Gamma 풀블리드 배경 사용)")
    elif html_slides_enabled and not bridge_available:
        # 명시적 안내만 추가 — 후속 mermaid/matplotlib는 그대로 진행
        pass

    # === Tier 0.6 — Vertex AI image models (Nano Banana Pro / Imagen 4) =====
    # OPT-IN via AE_ENABLE_VERTEX_IMAGE=1 + GOOGLE_APPLICATION_CREDENTIALS path.
    # Bedrock에 동급 텍스트-정확 이미지 모델이 없는 상황에서 사용자 결정으로
    # Vertex AI 호출을 예외 허용. LLM/추론은 그대로 Bedrock Gateway 경유 —
    # Vertex 호출은 *이미지 생성 한정*.
    #
    # Strengths vs Bedrock Stability:
    #   - Renders text accurately (Korean + English) — no fake glyphs
    #   - Composes 14+ objects coherently per scene
    #   - Photo-realistic at 4K
    #
    # 호출 흐름:
    #   for each visual section without HTML background → Nano Banana Pro
    #     → if success: section_diagrams[idx] = path
    #     → if fail: fall through to mermaid/matplotlib (no degradation)
    vertex_image_enabled = (
        _vertex_on
        and visual_intent
        and any(e in needed_exts for e in ("pptx", "pdf", "docx"))
        and os.environ.get("AE_PREFER_EDITABLE_DIAGRAM", "1") == "0"  # 편집 가능 우선이면 통짜 이미지 OFF
    )
    if vertex_image_enabled:
        try:
            from ai_engine.vertex_image_module import get_vertex_image_client
            _vertex_client = get_vertex_image_client(aws_profile=aws_profile or "")
        except Exception as _e_vimp:
            print(f"[ForceGenerate] Vertex import failed: {_e_vimp}")
            _vertex_client = None

        if _vertex_client and _vertex_client.enabled:
            # 비용/레이턴시 한계: 첫 8개 섹션만 시도. HTML 브릿지로 이미 커버된
            # 섹션은 제외. AE_VERTEX_MODEL_CLASS 로 capability 변경 가능
            # (default: image_generation_high_quality → Nano Banana Pro).
            _vertex_class = os.environ.get(
                "AE_VERTEX_MODEL_CLASS", "image_generation_high_quality"
            ).strip() or "image_generation_high_quality"
            vertex_section_inputs = []
            for _idx, _sec in enumerate(sections[:8]):
                if _idx in section_backgrounds:
                    continue
                _heading = (_sec.get("heading") or "").strip()
                _body = (_sec.get("body") or "").strip()
                if not (_heading or _body):
                    continue
                vertex_section_inputs.append((_idx, _heading, _body))

            async def _try_vertex_for_section(_idx_v, _heading_v, _body_v):
                # Build a slide-aware prompt — Nano Banana Pro CAN render Korean
                # text correctly, so we can safely include the heading literally.
                # But we still steer toward "professional infographic" style and
                # forbid watermarks / fake logos.
                _topic = (_heading_v or _body_v[:80] or "professional concept")[:120]
                # Active_Template + Vertex 활성 시 Style_Profile의 색/폰트 토큰을
                # 프롬프트 끝에 데이터로만 주입한다 (요구사항 7.4). vertex_image_module.py
                # 시그니처는 변경하지 않고 호출부에서 프롬프트 문자열만 조립한다.
                # _tpl_profile은 _force_generate_from_text 스코프의 로컬 변수로,
                # 본 클로저에서 직접 참조한다. None(무템플릿)이면 _style_hint는 빈
                # 문자열이라 기존 프롬프트가 바이트 단위로 보존된다(요구사항 7.5).
                _style_hint = ""
                if _tpl_profile:  # Active_Template 지정 + Vertex 활성 (요구사항 7.4)
                    _primary = _tpl_profile.get("primaryColor")
                    _accent = _tpl_profile.get("accentColor")
                    _hfont = _tpl_profile.get("headingFont")
                    _bfont = _tpl_profile.get("bodyFont")
                    _parts = []
                    # 토큰이 부재/무효(falsy)이면 그 토큰만 힌트에서 제외하고 나머지는
                    # 포함하며 렌더링은 중단 없이 계속한다 (요구사항 7.6).
                    if _primary and _accent:
                        _parts.append(f"Color palette: primary {_primary}, accent {_accent}.")
                    elif _primary:
                        _parts.append(f"Color palette: primary {_primary}.")
                    elif _accent:
                        _parts.append(f"Color palette: accent {_accent}.")
                    if _hfont and _bfont:
                        _parts.append(
                            f"Typography style cues: heading font \"{_hfont}\", "
                            f"body font \"{_bfont}\"."
                        )
                    elif _hfont:
                        _parts.append(f"Typography style cues: heading font \"{_hfont}\".")
                    elif _bfont:
                        _parts.append(f"Typography style cues: body font \"{_bfont}\".")
                    if _parts:
                        _style_hint = " " + " ".join(_parts)
                _prompt = (
                    f"Professional business infographic illustration for the slide "
                    f"titled \"{_topic}\". Modern flat design, clean layout, "
                    f"premium vector line iconography with consistent stroke weight, neutral color palette with subtle "
                    f"accent gradients. The image should visually summarize the "
                    f"concept and be safe to use in a corporate presentation. "
                    f"NO emoji, NO childish clipart, NO cartoon stickers."
                    f"{_style_hint}"
                )
                _negative = "watermark, fake logo, brand name, distorted UI, unreadable artifacts, emoji, emoji-style icons, childish clipart, cartoon stickers"
                try:
                    _vres = await _vertex_client.generate(
                        prompt=_prompt,
                        model_class=_vertex_class,
                        aspect_ratio="16:9",
                        negative_prompt=_negative,
                        timeout=60,
                    )
                    if "error" in _vres:
                        print(
                            f"[ForceGenerate] section[{_idx_v}] Vertex error: "
                            f"{_vres.get('error')} {(_vres.get('detail') or '')[:120]}"
                        )
                        return _idx_v, ""
                    _images = _vres.get("images") or []
                    if not _images:
                        return _idx_v, ""
                    _img_b64 = _images[0]
                    # Save to .generated/ in the same convention as
                    # _tool_generate_image (path returned is relative).
                    import time as _t_v, hashlib as _h_v, base64 as _b_v
                    _local_root_v = _resolve_local_root(project_path)
                    _gen_dir_v = os.path.join(_local_root_v, ".generated")
                    os.makedirs(_gen_dir_v, exist_ok=True)
                    _ts_v = str(int(_t_v.time() * 1000))
                    _short_v = _h_v.md5(_prompt.encode()).hexdigest()[:6]
                    _model_v = (_vres.get("model") or "vertex").split("/")[-1].replace(
                        "-", "")[:24]
                    _filename_v = f"vertex-{_ts_v}-{_short_v}-{_model_v}-{_idx_v}.png"
                    _abs_v = os.path.join(_gen_dir_v, _filename_v)
                    try:
                        with open(_abs_v, "wb") as _f_v:
                            _f_v.write(_b_v.b64decode(_img_b64))
                    except Exception as _e_save:
                        print(f"[ForceGenerate] Vertex save failed: {_e_save}")
                        return _idx_v, ""
                    return _idx_v, f".generated/{_filename_v}"
                except Exception as _e_v:
                    print(f"[ForceGenerate] section[{_idx_v}] Vertex exception: {_e_v}")
                    return _idx_v, ""

            if vertex_section_inputs:
                print(
                    f"[ForceGenerate] Vertex tier 0.6 활성 — "
                    f"sections={len(vertex_section_inputs)}, model_class={_vertex_class}"
                )
                try:
                    _vertex_results = await asyncio.gather(
                        *(
                            _try_vertex_for_section(_i, _h, _b)
                            for (_i, _h, _b) in vertex_section_inputs
                        ),
                        return_exceptions=True,
                    )
                    _vok = 0
                    for _r in _vertex_results:
                        if isinstance(_r, Exception):
                            print(f"[ForceGenerate] Vertex task exception: {_r}")
                            continue
                        _idx_r, _rel_r = _r
                        if _rel_r:
                            section_diagrams[_idx_r] = _rel_r
                            if not shared_native_image:
                                shared_native_image = _rel_r
                            _vok += 1
                    print(
                        f"[ForceGenerate] Vertex images: "
                        f"{_vok}/{len(vertex_section_inputs)} sections got "
                        f"Nano Banana / Imagen PNGs"
                    )
                except Exception as _e_vg:
                    print(f"[ForceGenerate] Vertex gather 실패: {_e_vg}")
        else:
            print(
                "[ForceGenerate] Vertex tier 0.6 — flag set but client disabled "
                "(check GOOGLE_APPLICATION_CREDENTIALS path or google-auth install)"
            )

    # === Tier 0.5 — Bedrock image models (Stability / Nova Canvas / Titan) ===
    # OFF BY DEFAULT — generative image models cannot reliably render text
    # (Korean OR English): they invent fake glyphs (킹딕록켠, TEUQUB, PEDUCTKIN)
    # that misalign with slide content. Even with `NO TEXT, NO LABELS` prompts,
    # the output is visually noisy and unrelated to the section body.
    #
    # Diagram-style slide bodies now route to mermaid (LLM-generated, Korean-
    # accurate) and python-pptx native renderers. Generative image models are
    # reserved for the cover/hero photo via Tier 0.7 below.
    #
    # To force the old behavior for testing/debugging, set
    # AE_ENABLE_BEDROCK_SLIDE_IMAGES=1.
    bedrock_image_enabled = (
        os.environ.get("AE_ENABLE_BEDROCK_SLIDE_IMAGES", "").strip() == "1"
        and visual_intent
        and not _image_gen_is_circuit_broken()
        and any(e in needed_exts for e in ("pptx", "pdf", "docx"))
    )
    if bedrock_image_enabled:
        # 비용/레이턴시 한계: 첫 8개 섹션만 시도. HTML 브릿지로 이미 커버된
        # 섹션은 제외 (전면 배경이 더 높은 퀄리티).
        bedrock_section_inputs = []
        for _idx, _sec in enumerate(sections[:8]):
            if _idx in section_backgrounds:
                continue
            _heading = (_sec.get("heading") or "").strip()
            _body = (_sec.get("body") or "").strip()
            if not (_heading or _body):
                continue
            # Structural sections (paths, arrow chains, markdown tables) go to
            # mermaid for accurate Korean labels — Bedrock cannot render Hangul.
            # _looks_structural is already imported in this file.
            if _looks_structural(_body, _heading, _body):
                print(
                    f"[ForceGenerate] section[{_idx}] structural → "
                    f"skip Bedrock, route to mermaid"
                )
                continue
            bedrock_section_inputs.append((_idx, _heading, _body))

        async def _try_bedrock_for_section(_idx_b, _heading_b, _body_b):
            _prompt = _build_section_image_prompt(
                _heading_b, _body_b, description or ""
            )
            try:
                _raw = await _tool_generate_image(
                    {"prompt": _prompt, "size": "1024x1024"},
                    project_path,
                    aws_profile=aws_profile,
                    bedrock_user=bedrock_user,
                )
                _parsed = {}
                try:
                    _parsed = json.loads(_raw)
                except (json.JSONDecodeError, TypeError):
                    _parsed = {}
                if (
                    isinstance(_parsed, dict)
                    and "error" not in _parsed
                    and _parsed.get("path")
                ):
                    return _idx_b, _parsed["path"]
            except Exception as _e_b:
                print(
                    f"[ForceGenerate] section[{_idx_b}] Bedrock image 예외: {_e_b}"
                )
            return _idx_b, ""

        if bedrock_section_inputs:
            print(
                f"[ForceGenerate] Bedrock image tier 0.5 활성 — "
                f"sections={len(bedrock_section_inputs)}"
            )
            try:
                bedrock_results = await asyncio.gather(
                    *(
                        _try_bedrock_for_section(_i, _h, _b)
                        for (_i, _h, _b) in bedrock_section_inputs
                    ),
                    return_exceptions=True,
                )
                _ok_count = 0
                for _r in bedrock_results:
                    if isinstance(_r, Exception):
                        print(
                            f"[ForceGenerate] Bedrock image task exception: {_r}"
                        )
                        continue
                    _idx_r, _rel_r = _r
                    if _rel_r:
                        section_diagrams[_idx_r] = _rel_r
                        if not shared_native_image:
                            shared_native_image = _rel_r
                        _ok_count += 1
                print(
                    f"[ForceGenerate] Bedrock images: "
                    f"{_ok_count}/{len(bedrock_section_inputs)} sections got "
                    f"Stability/Nova/Titan PNGs"
                )
            except Exception as _e_g:
                print(f"[ForceGenerate] Bedrock image gather 실패: {_e_g}")

    # === Tier 0.7 — Hero/cover photo via Bedrock (optional) ==================
    # Stability Ultra produces excellent abstract photographic backgrounds when
    # the prompt is photographic (no text expected). Use ONE such image for the
    # cover slide only — never for diagram bodies. Opt-in via
    # AE_BEDROCK_HERO_IMAGE=1.
    #
    # Constraint compliance: routed through _tool_generate_image which goes
    # through _get_gw().invoke_model (Bedrock Gateway) — gateway.md rules.
    hero_image_enabled = (
        os.environ.get("AE_BEDROCK_HERO_IMAGE", "").strip() == "1"
        and visual_intent
        and not _image_gen_is_circuit_broken()
        and any(e in needed_exts for e in ("pptx", "pdf"))
        and not section_backgrounds  # HTML hero already covered
    )
    hero_image_path = ""
    _shared_body_bg = ""
    if hero_image_enabled:
        # Photographic prompt — NO text, NO labels, NO diagram words.
        hero_topic = (title or description or "professional concept")[:80]
        # Strip Korean to avoid the model attempting glyph synthesis.
        import re as _re_h
        hero_ascii = _re_h.sub(r"[^\x00-\x7F]+", " ", hero_topic)
        hero_ascii = _re_h.sub(r"\s+", " ", hero_ascii).strip() or "professional concept"
        hero_prompt = (
            f"professional cinematic photograph: {hero_ascii}, "
            "abstract composition, soft natural lighting, depth of field, "
            "neutral color palette, high resolution, "
            "NO TEXT, NO LABELS, NO LETTERS, NO LOGO"
        )[:250]
        try:
            hero_raw = await _tool_generate_image(
                {"prompt": hero_prompt, "size": "1536x1024", "style": "photographic"},
                project_path, aws_profile=aws_profile, bedrock_user=bedrock_user,
            )
            hero_parsed = json.loads(hero_raw) if hero_raw else {}
            if isinstance(hero_parsed, dict) and "error" not in hero_parsed and hero_parsed.get("path"):
                hero_image_path = hero_parsed["path"]
                print(f"[ForceGenerate] Hero photo: {hero_image_path}")
        except Exception as _e_hero:
            print(f"[ForceGenerate] Hero photo failed: {_e_hero}")

    # === Tier 0.65 — Vertex 고품질 표지 배경(기본 ON, Vertex 활성 시) ===========
    # 사용자 요구: "Nano Banana Pro급 이미지 + 편집 가능". 표지에 Vertex 고품질
    # 배경을 깔고, 그 위에 편집 가능한 네이티브 제목/KPI(스크림으로 가독성)를 올린다.
    # 본문은 편집 가능 네이티브 다이어그램 유지(편집성 우선). AE_DISABLE_VERTEX_HERO=1로 끔.
    if (not hero_image_path and _vertex_on and visual_intent
            and any(e in needed_exts for e in ("pptx", "pdf"))
            and not section_backgrounds
            and os.environ.get("AE_DISABLE_VERTEX_HERO", "").strip() != "1"):
        try:
            _vh_topic = (title or description or "professional concept")[:90]
            _vh_prompt = (
                "Professional abstract background illustration for a presentation cover about: "
                f"{_vh_topic}. Modern corporate style, smooth gradient, subtle geometric depth, "
                "clean empty negative space on the left third for a title, cinematic soft lighting, "
                "no text, no words, no letters, no captions, high resolution, 16:9."
            )
            _vh_raw = await _tool_generate_image(
                {"prompt": _vh_prompt, "size": "1280x720", "style": "digital art"},
                project_path, aws_profile=aws_profile, bedrock_user=bedrock_user,
            )
            _vh_parsed = json.loads(_vh_raw) if _vh_raw else {}
            if (isinstance(_vh_parsed, dict) and "error" not in _vh_parsed
                    and _vh_parsed.get("path")):
                hero_image_path = _vh_parsed["path"]
                print(f"[ForceGenerate] Vertex 표지 배경: {hero_image_path} (via={_vh_parsed.get('via')})")
            else:
                print(f"[ForceGenerate] Vertex 표지 배경 미생성: {str(_vh_parsed)[:160]}")
        except Exception as _e_vh:
            print(f"[ForceGenerate] Vertex 표지 배경 실패: {str(_e_vh)[:160]}")

    # === Tier 0.66 — Vertex 본문 공유 배경(기본 ON) =========================
    # 모든 본문 슬라이드에 동일한 얙은 고품질 배경을 깔아 표지와 통일된 고급감을
    # 준다. 그 위에는 흠 콘텐츠 카드 + 편집 가능 다이어그램/텍스트가 올라온다.
    # AE_DISABLE_VERTEX_BODY_BG=1로 끕.
    if (_vertex_on and visual_intent
            and any(e in needed_exts for e in ("pptx",))
            and not section_backgrounds
            and os.environ.get("AE_DISABLE_VERTEX_BODY_BG", "").strip() != "1"):
        try:
            _bb_prompt = (
                "Very light, pale, minimal abstract background for presentation content slides. "
                "Soft white and faint pastel gradient, lots of bright white negative space, "
                "subtle thin geometric lines in one corner, clean, airy, no dark areas, "
                "no text, no words, no letters, high resolution, 16:9."
            )
            _bb_raw = await _tool_generate_image(
                {"prompt": _bb_prompt, "size": "1280x720", "style": "digital art"},
                project_path, aws_profile=aws_profile, bedrock_user=bedrock_user,
            )
            _bb_parsed = json.loads(_bb_raw) if _bb_raw else {}
            if (isinstance(_bb_parsed, dict) and "error" not in _bb_parsed
                    and _bb_parsed.get("path")):
                _shared_body_bg = _bb_parsed["path"]
                print(f"[ForceGenerate] Vertex 본문 배경: {_shared_body_bg}")
        except Exception as _e_bb:
            print(f"[ForceGenerate] Vertex 본문 배경 실패: {str(_e_bb)[:160]}")

    if use_native:
        # === Mermaid 우선 전략 ===
        # 1) 섹션별로 LLM에게 mermaid 코드 생성 → mermaid.ink로 PNG 렌더 (병렬).
        # 2) Mermaid 실패한 섹션만 matplotlib classifier fallback.
        # AE_DISABLE_MERMAID=1 환경변수로 mermaid 끄고 matplotlib만 쓸 수 있음.
        mermaid_enabled = os.environ.get("AE_DISABLE_MERMAID", "").strip() != "1" \
            and not _image_gen_is_circuit_broken()  # 회로 차단 시 외부 HTTP도 위험할 수 있음

        diagram_cache: dict = {}
        gw_for_mermaid = None
        mermaid_model = ""
        if mermaid_enabled:
            try:
                gw_for_mermaid = _get_gw(aws_profile, bedrock_user)
                # Sonnet이 mermaid 정확도 가장 높음 + 도구 호출 부담 없으니 Sonnet 써도 됨
                mermaid_model = _specialized_model_for_task(
                    "file_generation", "",
                    aws_profile=aws_profile, bedrock_user=bedrock_user,
                )
            except Exception as e:
                print(f"[ForceGenerate] mermaid 환경 준비 실패: {e}")
                mermaid_enabled = False

        # 1단계 — 섹션 의미 분류로 어떤 섹션이 다이어그램 필요한지 결정
        section_specs = []  # [(idx, heading, body, kind, content)]
        # section_diagram_specs는 use_native 블록 밖(상단)에서 이미 초기화됨.
        # 여기서는 채우기만 한다(재초기화 금지 — 위 주석 참고).
        for idx, sec in enumerate(sections):
            heading = (sec.get("heading") or "").strip()
            body = (sec.get("body") or "").strip()
            kind, content_for_diagram = _classify_section_diagram(heading, body, title or "")
            if not kind:
                continue
            section_specs.append((idx, heading, body, kind, content_for_diagram))
            section_diagram_specs[idx] = {"type": kind, "content": content_for_diagram}

        # 2단계 — Mermaid 병렬 생성 (LLM 호출 + mermaid.ink 렌더)
        async def _try_mermaid_for_section(idx, heading, body):
            cache_key = ("mermaid", (heading + body[:300])[:300])
            if cache_key in diagram_cache:
                return idx, diagram_cache[cache_key]
            try:
                code = await _llm_generate_mermaid(
                    gw_for_mermaid, mermaid_model,
                    section_heading=heading, section_body=body,
                    doc_context=(title or description or "")[:200],
                    style_profile=_tpl_profile,
                )
                if not code:
                    return idx, ""
                render_json = await _render_mermaid_to_png(code, project_path=project_path, timeout=30)
                parsed = {}
                try:
                    parsed = json.loads(render_json)
                except (json.JSONDecodeError, TypeError):
                    pass
                if isinstance(parsed, dict) and parsed.get("path"):
                    diagram_cache[cache_key] = parsed["path"]
                    return idx, parsed["path"]
                print(f"[ForceGenerate] section[{idx}] mermaid 렌더 실패: {render_json[:200]}")
                return idx, ""
            except Exception as e:
                print(f"[ForceGenerate] section[{idx}] mermaid 예외: {e}")
                return idx, ""

        if mermaid_enabled and section_specs:
            # Mermaid runs for every visual-intent section that's not covered by a
            # higher-tier background. This is the primary diagram path now that
            # Bedrock image generation is reserved for hero photos only — every
            # visual-intent section gets a Korean-accurate mermaid diagram (or
            # matplotlib fallback below) instead of a fake-glyph generative PNG.
            tasks = [
                _try_mermaid_for_section(idx, heading, body)
                for (idx, heading, body, _kind, _content) in section_specs
                # Skip mermaid for sections that already have an HTML background —
                # the HTML slide is higher quality and we don't want a duplicate
                # diagram on top of the full-bleed background.
                # Also skip sections already populated by tier 0.5 Bedrock image
                # models (when AE_ENABLE_BEDROCK_SLIDE_IMAGES=1) — overwriting a
                # Stability/Nova-Canvas PNG with a mermaid render would defeat the
                # opt-in. With the default config (Tier 0.5 OFF), section_diagrams
                # is empty here so mermaid fans out to every visual section.
                if idx not in section_backgrounds and idx not in section_diagrams
            ]
            mermaid_results = await asyncio.gather(*tasks, return_exceptions=True)
            for r in mermaid_results:
                if isinstance(r, Exception):
                    continue
                idx, rel = r
                if rel:
                    section_diagrams[idx] = rel
                    if not shared_native_image:
                        shared_native_image = rel
                    print(f"[ForceGenerate] section[{idx}] mermaid OK → {rel}")

        # 3단계 — Mermaid 실패한 섹션만 matplotlib fallback
        for (idx, heading, body, kind, content_for_diagram) in section_specs:
            if idx in section_diagrams:
                continue
            cache_key = (kind, content_for_diagram[:200])
            if cache_key in diagram_cache:
                section_diagrams[idx] = diagram_cache[cache_key]
                continue
            try:
                diagram_result = await _tool_generate_native_diagram(
                    diagram_type=kind,
                    title=(heading or title or "Diagram")[:80],
                    content=content_for_diagram,
                    project_path=project_path,
                    palette=_tpl_palette,
                )
                parsed = {}
                try:
                    parsed = json.loads(diagram_result)
                except (json.JSONDecodeError, TypeError):
                    pass
                if isinstance(parsed, dict) and parsed.get("path"):
                    section_diagrams[idx] = parsed["path"]
                    diagram_cache[cache_key] = parsed["path"]
                    if not shared_native_image:
                        shared_native_image = parsed["path"]
                    print(f"[ForceGenerate] section[{idx}] matplotlib fallback ({kind}) → {parsed['path']}")
            except Exception as e:
                print(f"[ForceGenerate] section[{idx}] matplotlib fallback 예외: {e}")

        # 4단계 — 각 섹션마다 matplotlib을 보장 fallback으로.
        # mermaid/Bedrock가 채우지 못한 섹션을 텍스트만 남기지 않도록,
        # heading/body 기반으로 section-specific 다이어그램을 한 장씩 그려서
        # section_diagrams[idx]에 채운다. 이렇게 하면 imagePrompt 분기에 떨어지지 않음.
        for (idx, heading, body, kind, content_for_diagram) in section_specs:
            if idx in section_diagrams or idx in section_backgrounds:
                continue
            try:
                fallback_kind = kind or _classify_section_diagram(heading, body, title or "")[0] or "block"
                # 섹션 콘텐츠가 비어있으면 heading만이라도 사용 — 빈 다이어그램은 만들지 않음
                content_for_section = (content_for_diagram or body or heading or "").strip()
                if not content_for_section:
                    continue
                diagram_result = await _tool_generate_native_diagram(
                    diagram_type=fallback_kind,
                    title=(heading or title or "Diagram")[:80],
                    content=content_for_section[:4000],
                    project_path=project_path,
                    palette=_tpl_palette,
                )
                parsed_sec = {}
                try:
                    parsed_sec = json.loads(diagram_result)
                except (json.JSONDecodeError, TypeError):
                    pass
                if isinstance(parsed_sec, dict) and parsed_sec.get("path"):
                    section_diagrams[idx] = parsed_sec["path"]
                    if not shared_native_image:
                        shared_native_image = parsed_sec["path"]
                    print(
                        f"[ForceGenerate] section[{idx}] matplotlib guaranteed fallback "
                        f"({fallback_kind}) → {parsed_sec['path']}"
                    )
            except Exception as e:
                print(f"[ForceGenerate] section[{idx}] guaranteed fallback 예외: {e}")

        # 4-bis: 그래도 shared_native_image가 비어있으면 (모든 섹션이 빈 콘텐츠) 문서 전체 한 장
        if not shared_native_image:
            try:
                fallback_kind = _classify_section_diagram(
                    "", final_text or description or "", title or ""
                )[0] or "block"
                diagram_content_full = (final_text or description or title or "").strip()
                if diagram_content_full:
                    diagram_result = await _tool_generate_native_diagram(
                        diagram_type=fallback_kind,
                        title=(title or description or "Diagram")[:80],
                        content=diagram_content_full,
                        project_path=project_path,
                        palette=_tpl_palette,
                    )
                    parsed_doc = {}
                    try:
                        parsed_doc = json.loads(diagram_result)
                    except (json.JSONDecodeError, TypeError):
                        pass
                    if isinstance(parsed_doc, dict) and parsed_doc.get("path"):
                        shared_native_image = parsed_doc["path"]
                        print(
                            f"[ForceGenerate] document-wide fallback diagram: "
                            f"{fallback_kind} → {shared_native_image}"
                        )
            except Exception as e:
                print(f"[ForceGenerate] document fallback 예외: {e}")

    out = []
    # 검증 시 도구가 저장한 실제 위치를 따라가야 함 — _resolve_relative_for_verify는
    # project_path → AE_GENERATED_ROOT → ~/.agentic-editor → tempdir 순으로 후보를
    # 시도해 첫 존재 경로를 반환. 도구는 _resolve_local_root와 같은 우선순위를 사용.
    project_root = project_path if (project_path and os.path.isdir(project_path)) else os.getcwd()

    for ext in needed_exts:
        try:
            if ext == "pdf":
                section_dicts = []
                # 섹션별 의미 인식 다이어그램 — section_diagrams[idx]가 있으면 그것을, 없으면
                # shared_native_image fallback. 첫 섹션엔 무조건 이미지 부여 (원본 버그: 헤더 빈
                # 섹션에서 imageFile 안 들어가 PDF가 24KB로 끝나던 문제 회귀 방지).
                _img_attached = False
                for s_idx, s in enumerate(sections):
                    sd = {"heading": s["heading"], "body": s["body"]}
                    # Tier 0: HTML slide background (capped to 12cm in PDF so text still fits).
                    bg_for_this = section_backgrounds.get(s_idx)
                    if bg_for_this:
                        sd["slideBackground"] = bg_for_this
                        _img_attached = True
                    else:
                        img_for_this = section_diagrams.get(s_idx) or (shared_native_image if not _img_attached else "")
                        if visual_intent and img_for_this:
                            sd["imageFile"] = img_for_this
                            _img_attached = True
                        elif (
                            _bedrock_slide_on
                            and visual_intent
                            and not _img_attached
                            and not _image_gen_is_circuit_broken()
                        ):
                            # Only when AE_ENABLE_BEDROCK_SLIDE_IMAGES=1 — otherwise
                            # let _tool_generate_pdf produce a text-only section.
                            sd["imagePrompt"] = _build_section_image_prompt(
                                s["heading"], s["body"], description or ""
                            )
                            _img_attached = True
                    section_dicts.append(sd)
                inp = {
                    "title": title or "Document",
                    "sections": section_dicts,
                }
                tout = await _tool_generate_pdf(inp, project_path,
                    aws_profile=aws_profile, bedrock_user=bedrock_user)
            elif ext == "pptx":
                slide_dicts = []
                _img_attached_pptx = False
                for s_idx, s in enumerate(sections[:10]):
                    sd = {
                        "title": s["heading"][:80] or (title or "Slide")[:80],
                        "bullets": _extract_bullets(s["body"])[:6] or [s["body"][:200]],
                    }
                    # Tier 0: HTML slide background (full-bleed) — best quality.
                    # 이슈 4 — 옵션 A: 풀블리드 배경은 편집 불가 통짜 슬라이드를 만든다.
                    # AE_ENABLE_HTML_SLIDES=1 옵트인일 때만 배경으로 사용하고,
                    # 기본값에서는 보조 이미지(imageFile, 우측 절반)로 강등해 텍스트
                    # placeholder가 항상 편집 가능하도록 한다.
                    bg_for_this = section_backgrounds.get(s_idx) if html_slides_enabled else None
                    if bg_for_this:
                        sd["slideBackground"] = bg_for_this
                        _img_attached_pptx = True
                    else:
                        # Tier 1/2: mermaid/matplotlib diagram (right-half embed).
                        # HTML 옵트인이 아니어도 section_backgrounds에 캡쳐가 있으면
                        # (과거 캐시 등) 보조 이미지로 활용 — 편집 가능 유지.
                        img_for_this = (
                            section_diagrams.get(s_idx)
                            or (section_backgrounds.get(s_idx) if not html_slides_enabled else None)
                            or (shared_native_image if not _img_attached_pptx else "")
                        )
                        # 문제2 — 네이티브(편집 가능) 다이어그램 우선.
                        # 이 섹션에 다이어그램 spec(type+content)이 있으면, 통짜 PNG 대신
                        # python-pptx 네이티브 도형으로 그리도록 슬라이드 dict에 spec을 전달한다.
                        # _tool_generate_pptx가 nativeDiagram을 받으면 imageFile보다 우선해
                        # 도형/텍스트/커넥터를 직접 조립한다(다운로드 후 개별 편집 가능).
                        # AE_DISABLE_NATIVE_DIAGRAM=1이면 기존 PNG 경로로 폴백.
                        _native_off = os.environ.get("AE_DISABLE_NATIVE_DIAGRAM", "") == "1"
                        _editable_pref2 = os.environ.get("AE_PREFER_EDITABLE_DIAGRAM", "1") != "0"
                        _diag_spec = section_diagram_specs.get(s_idx)
                        if _editable_pref2 and not _native_off:
                            # 편집 가능 우선: 통짜 PNG/저품질 휴리스틱 spec을 넣지 않는다.
                            # 깨끗한 title+bullets만 넘기면 _tool_generate_pptx의 LLM 구조화가
                            # 이 섹션을 고품질 편집 가능 네이티브 다이어그램으로 변환한다.
                            pass
                        elif visual_intent and _diag_spec and not _native_off:
                            sd["nativeDiagram"] = _diag_spec
                            if img_for_this:
                                sd["imageFile"] = img_for_this
                            _img_attached_pptx = True
                        elif visual_intent and img_for_this:
                            sd["imageFile"] = img_for_this
                            _img_attached_pptx = True
                        elif (
                            _bedrock_slide_on
                            and visual_intent
                            and not _img_attached_pptx
                            and not _image_gen_is_circuit_broken()
                        ):
                            # Only when AE_ENABLE_BEDROCK_SLIDE_IMAGES=1 — otherwise
                            # let _tool_generate_pptx produce a text-only slide.
                            sd["imagePrompt"] = _build_section_image_prompt(
                                s["heading"], s["body"], description or ""
                            )
                            _img_attached_pptx = True
                    # 본문도 표지급 고품질: 옅은 Vertex 공유 배경을 깔고
                    # _tool_generate_pptx가 그 위에 흰 카드+편집 가능 다이어그램을 그린다.
                    if (_shared_body_bg and visual_intent
                            and not sd.get("slideBackground")):
                        sd["slideBackground"] = _shared_body_bg
                    slide_dicts.append(sd)
                inp = {
                    "title": title or "Presentation",
                    "slides": slide_dicts or [{"title": title or "Slide 1", "bullets": [base_text[:200]]}],
                }
                # Use the first available background as the cover background too — gives
                # the title slide a polished look instead of plain white default layout.
                # 이슈 4 — 옵션 A: 표지 배경도 HTML 옵트인일 때만 사용 (표지 텍스트는
                # placeholder로 배경 위에 남아 편집 가능하지만, 기본은 깔끔한 네이티브 표지).
                first_bg = (next(iter(section_backgrounds.values()), "") if (section_backgrounds and html_slides_enabled) else "")
                if first_bg:
                    inp["coverBackground"] = first_bg
                elif hero_image_path:
                    # Hero photo (Tier 0.7) is the cover background when no HTML
                    # bridge available — single Bedrock call, photographic prompt,
                    # never used for diagram bodies.
                    inp["coverBackground"] = hero_image_path
                # 활성 템플릿이 유효할 때만 templatePath/templateId/styleProfile 주입 (요구사항 5.1, 5.3).
                # 무템플릿일 때는 이 키들을 추가하지 않아 inp가 기존(baseline)과 바이트 동일하다(요구사항 5.2).
                if _tpl_used:
                    inp["templatePath"] = _tpl_path
                    inp["templateId"] = template_id
                    inp["styleProfile"] = _tpl_profile
                tout = await _tool_generate_pptx(inp, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user)
            elif ext == "xlsx":
                # 마크다운 표가 있으면 그것을 사용, 없으면 섹션 헤더 + 본문 요약
                rows = _extract_md_table(base_text) or [
                    [s["heading"][:60], s["body"][:300]] for s in sections
                ]
                headers = ["Heading", "Content"] if not _extract_md_table(base_text) else None
                if headers is None:
                    # 첫 행이 헤더로 추출됨
                    md_table = _extract_md_table(base_text)
                    headers = md_table[0] if md_table else ["Col1", "Col2"]
                    rows = md_table[1:] if md_table and len(md_table) > 1 else []
                inp = {
                    "title": title or "Workbook",
                    "sheets": [{"name": "Summary", "headers": headers, "rows": rows or [["(empty)", "(empty)"]]}],
                }
                tout = await _tool_generate_xlsx(inp, project_path)
            elif ext == "docx":
                section_dicts = []
                _img_attached_docx = False
                for s_idx, s in enumerate(sections):
                    sd = {
                        "heading": s["heading"],
                        "level": 2,
                        "body": s["body"],
                        "bullets": _extract_bullets(s["body"]),
                    }
                    bg_for_this = section_backgrounds.get(s_idx)
                    if bg_for_this:
                        # DOCX uses imageFile path for backgrounds (capped width).
                        sd["imageFile"] = bg_for_this
                        _img_attached_docx = True
                    else:
                        img_for_this = section_diagrams.get(s_idx) or (shared_native_image if not _img_attached_docx else "")
                        if visual_intent and img_for_this:
                            sd["imageFile"] = img_for_this
                            _img_attached_docx = True
                        elif (
                            _bedrock_slide_on
                            and visual_intent
                            and not _img_attached_docx
                            and not _image_gen_is_circuit_broken()
                        ):
                            # Only when AE_ENABLE_BEDROCK_SLIDE_IMAGES=1 — otherwise
                            # let _tool_generate_docx produce a text-only section.
                            sd["imagePrompt"] = _build_section_image_prompt(
                                s["heading"], s["body"], description or ""
                            )
                            _img_attached_docx = True
                    section_dicts.append(sd)
                inp = {
                    "title": title or "Document",
                    "sections": section_dicts,
                }
                tout = await _tool_generate_docx(inp, project_path,
                    aws_profile=aws_profile, bedrock_user=bedrock_user)
            elif ext in ("png", "jpg", "jpeg"):
                # [native-diagram] 회로 차단 또는 구조형 콘텐츠면 matplotlib 우선
                if shared_native_image:
                    abs_native = os.path.join(_resolve_local_root(project_path), shared_native_image)
                    if os.path.isfile(abs_native):
                        tout = json.dumps({
                            "path": shared_native_image,
                            "model": "matplotlib (native)",
                            "sizeBytes": os.path.getsize(abs_native),
                        })
                    else:
                        tout = json.dumps({"error": "native-diagram-missing"})
                elif _image_gen_is_circuit_broken() or _looks_structural(description, title, final_text or ""):
                    d_type = "tree" if any(kw in (description or title or "").lower()
                                            for kw in ("폴더", "디렉토리", "tree", "구조")) else "flow"
                    diagram_content = (final_text or description or title or "")[:4000]
                    tout = await _tool_generate_native_diagram(
                        diagram_type=d_type,
                        title=(title or "Diagram")[:80],
                        content=diagram_content,
                        project_path=project_path,
                        palette=_tpl_palette,
                    )
                else:
                    inp = {"prompt": (description or title or "Architecture diagram")[:1500], "size": "1024x1024"}
                    tout = await _tool_generate_image(inp, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user)
            elif ext in ("md", "txt"):
                # write_file로 직접 작성
                fname = (target_files[0] if target_files else f".generated/{_safe_slug(title or 'note')}.{ext}")
                if not fname.startswith(".generated/") and not os.path.isabs(fname):
                    fname = ".generated/" + os.path.basename(fname)
                abs_path = fname if os.path.isabs(fname) else os.path.join(project_root, fname)
                os.makedirs(os.path.dirname(abs_path), exist_ok=True)
                with open(abs_path, "w", encoding="utf-8") as f:
                    f.write(base_text)
                tout = json.dumps({"path": fname, "model": "filesystem", "sizeBytes": os.path.getsize(abs_path)})
            else:
                continue

            try:
                parsed = json.loads(tout) if isinstance(tout, str) else tout
            except (json.JSONDecodeError, TypeError):
                parsed = None
            if not isinstance(parsed, dict) or "path" not in parsed or "error" in parsed:
                print(f"[ForceGenerate] {ext} 실패: {tout[:300]}")
                continue
            rel = parsed["path"]
            # 도구가 _resolve_local_root로 저장 → 검증도 같은 후보 우선순위로 lookup.
            # project_root에 없으면 AE_GENERATED_ROOT/~/.agentic-editor/tempdir도 시도.
            abs_path = _resolve_relative_for_verify(rel, project_path)
            if not os.path.isfile(abs_path):
                continue
            out.append((rel, {
                "path": rel,
                "absPath": abs_path,
                "size": os.path.getsize(abs_path),
                "tool": "generate_" + (ext if ext != "jpg" else "image"),
                # forced 경로지만 실제 사용된 라이브러리/엔진 이름이 와야 함.
                # 기본값은 deterministic-converter — system_fallback 라벨은 사용자가
                # 명시적으로 거부했으므로 절대 사용하지 않음.
                "model": parsed.get("model", "deterministic-converter"),
            }))
        except Exception as e:
            print(f"[ForceGenerate] {ext} 예외: {e}")
            continue

    return out


def _detect_visual_intent(user_prompt: str) -> bool:
    """사용자 프롬프트에서 시각 콘텐츠(이미지/다이어그램) 요구를 감지.

    KOR/EN 키워드 매칭으로 단순/결정적 판단. 어떤 키워드라도 등장하면 True.
    PDF/PPTX/DOCX 자동 imagePrompt 생성을 트리거하는 데 사용된다.
    """
    if not user_prompt:
        return False
    text = str(user_prompt).lower()
    keywords = (
        # Korean
        "시각", "이미지", "그림", "사진", "다이어그램", "차트", "그래프",
        "시각화", "일러스트", "도식", "그래픽",
        # Korean — 흐름/구조/발표 의도도 시각 슬라이드로 취급(고품질 tier 활성화)
        "흐름도", "흐름", "플로우", "프로세스", "분석", "아키텍처", "시스템",
        "구조", "구성도", "파이프라인", "슬라이드", "발표", "보고", "개요", "단계",
        # English
        "visual", "image", "picture", "photo", "diagram", "chart",
        "graph", "illustration", "infographic", "schema",
        "flow", "flowchart", "process", "architecture", "system", "pipeline",
        "workflow", "slide", "presentation", "overview", "structure",
    )
    return any(kw in text for kw in keywords)


def _build_section_image_prompt(heading: str, body: str, context: str = "") -> str:
    """섹션 → Stability/Titan 프롬프트 변환.

    핵심 제약:
      - Stability/Nova-Canvas/Titan 모델은 한글 글리프를 렌더링하지 못한다.
        Korean keywords are STRIPPED — only ASCII tokens survive.
      - 모든 프롬프트에 "NO TEXT / NO LABELS / NO LETTERS" 금지령 부착.
        라벨이 필요한 다이어그램은 mermaid가 담당 (한글 정확).
      - Bedrock은 "추상적 시각 비유"만 생성 — gears, network nodes,
        flowing shapes, layered planes (텍스트 없는 컨셉 일러스트).
    """
    import re as _re

    raw = f"{heading or ''} {body or ''} {context or ''}"
    # Strip Korean (Hangul Syllables + Jamo) and CJK ideographs.
    # Keeps only ASCII letters, digits, spaces, basic punct.
    ascii_only = _re.sub(r"[^\x00-\x7F]+", " ", raw)
    ascii_only = _re.sub(r"\s+", " ", ascii_only).strip().lower()

    is_flow = any(kw in ascii_only for kw in (
        "flow", "process", "workflow", "pipeline", "sequence", "step",
    ))
    is_arch = any(kw in ascii_only for kw in (
        "architecture", "system", "module", "component", "infrastructure",
        "service", "layer",
    ))
    is_data = any(kw in ascii_only for kw in (
        "data", "schema", "database", "table", "storage",
    ))
    is_ui = any(kw in ascii_only for kw in (
        "ui", "interface", "screen", "mockup", "wireframe", "dashboard",
    ))

    # NEGATIVE directives: prevent the model from inventing fake glyphs.
    no_text = (
        ", NO TEXT, NO LABELS, NO LETTERS, NO WORDS, NO TYPOGRAPHY"
        ", NO EMOJI, NO emoji-style icons, NO childish clipart, NO cartoon stickers"
        ", abstract conceptual illustration only, clean minimal symbolic shapes"
    )
    style = (
        "clean professional flat design, modern, neutral colors, soft gradients, "
        "premium vector line iconography with consistent stroke weight, "
        "enterprise presentation quality"
    )

    if is_flow:
        # Abstract flow visualization — flowing curves, no boxes-with-labels.
        base = (
            "abstract flow concept: smooth curving lines, flowing shapes, "
            "directional gradient, sense of motion and connection, isometric view"
        )
    elif is_arch:
        # Layered planes / nested structures, no labels.
        base = (
            "abstract system architecture: stacked translucent layers, "
            "interconnected geometric shapes, depth, isometric perspective"
        )
    elif is_data:
        # Data flow as light streams between abstract nodes.
        base = (
            "abstract data concept: glowing nodes connected by light streams, "
            "geometric pattern, depth, dark background"
        )
    elif is_ui:
        # Generic device silhouette without UI text.
        base = (
            "abstract device silhouette with empty geometric panels, "
            "minimalist, clean lines, no interface text"
        )
    else:
        base = (
            "abstract conceptual illustration: smooth geometric shapes, "
            "soft gradients, depth, modern minimal design"
        )

    prompt = f"{base}, {style}{no_text}"
    return prompt[:400]


def _split_into_sections(text: str, fallback_title: str = "본문") -> list:
    """마크다운 텍스트를 헤더 단위로 분리. 헤더가 없으면 단일 섹션."""
    if not text:
        return [{"heading": fallback_title, "body": ""}]
    lines = text.splitlines()
    sections = []
    current = {"heading": "", "body": []}
    for line in lines:
        m = re.match(r"^(#{1,6})\s+(.+)$", line.strip())
        if m:
            if current["heading"] or current["body"]:
                sections.append({"heading": current["heading"] or fallback_title,
                                 "body": "\n".join(current["body"]).strip()})
            current = {"heading": m.group(2).strip(), "body": []}
        else:
            current["body"].append(line)
    if current["heading"] or current["body"]:
        sections.append({"heading": current["heading"] or fallback_title,
                         "body": "\n".join(current["body"]).strip()})
    return sections or [{"heading": fallback_title, "body": text.strip()}]


def _extract_bullets(text: str) -> list:
    """텍스트에서 불릿 항목(- / * / 숫자.)을 추출."""
    if not text:
        return []
    bullets = []
    for line in text.splitlines():
        stripped = line.strip()
        m = re.match(r"^(?:[-*•]\s+|\d+\.\s+)(.+)$", stripped)
        if m:
            bullets.append(m.group(1).strip())
    return bullets


def _extract_md_table(text: str) -> list:
    """마크다운 표를 [[헤더...], [행1...], ...] 형태로 추출. 없으면 빈 리스트."""
    if not text:
        return []
    lines = text.splitlines()
    rows = []
    in_table = False
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("|") and stripped.endswith("|"):
            cells = [c.strip() for c in stripped.strip("|").split("|")]
            # 구분선 (---|---) 스킵
            if all(re.match(r"^:?-+:?$", c) for c in cells if c):
                in_table = True
                continue
            rows.append(cells)
            in_table = True
        elif in_table:
            break
    return rows


def _slug_from_title(s: str) -> str:
    """파일명용 슬러그 — 한국어 포함 보존, 특수문자만 제거.

    Korean/영숫자/한글/하이픈/언더스코어만 유지. 공백은 하이픈으로.
    """
    import re as _re
    if not s:
        return ""
    s = s.strip()
    # 영숫자, 한글, 일본어, 중국어, 하이픈, 언더스코어, 공백만 유지
    cleaned = _re.sub(r"[^\w\s\u3131-\u318E\uAC00-\uD7A3\u3040-\u309F\u30A0-\u30FF\u4E00-\u9FFF-]+", "", s, flags=_re.UNICODE)
    # 공백 → 하이픈
    cleaned = _re.sub(r"\s+", "-", cleaned)
    # 연속 하이픈 정리
    cleaned = _re.sub(r"-+", "-", cleaned).strip("-_")
    return cleaned[:30]


def _safe_slug(s: str) -> str:
    """파일명용 안전한 슬러그."""
    if not s:
        return "doc"
    s = re.sub(r"[^a-zA-Z0-9가-힣\-_\s]+", "", s)
    s = re.sub(r"\s+", "-", s.strip())
    return (s[:40] or "doc").lower()


async def _orchestrator_run_agent(
    gw, stream_model: str, subtask: dict, project_path: str,
    base_system_prompt: str, emit_queue: asyncio.Queue,
    aws_profile: str = "", bedrock_user: str = "", is_remote: bool = False,
    template_id: str = "",
):
    """하나의 하위 에이전트 실행 — 도구 루프 포함. 무한루프 방지를 위해 hard timeout."""
    # 환경 변수로 조정 가능 — 기본 30분, 도구 루프 50회 제한과 함께 동작
    agent_timeout = float(os.environ.get("AE_AGENT_TIMEOUT", "1800"))
    task_id = subtask.get("id", "?")
    role = subtask.get("role", "Worker")
    title = subtask.get("title", "")

    async def _run_inner():
        return await _orchestrator_run_agent_inner(
            gw, stream_model, subtask, project_path,
            base_system_prompt, emit_queue,
            aws_profile=aws_profile, bedrock_user=bedrock_user, is_remote=is_remote,
            template_id=template_id,
        )

    try:
        return await asyncio.wait_for(_run_inner(), timeout=agent_timeout)
    except asyncio.TimeoutError:
        msg = f"에이전트 시간 초과 ({int(agent_timeout)}초)"
        try:
            await emit_queue.put({"type": "agent_error", "taskId": task_id, "error": msg})
        except Exception:
            pass
        return {"taskId": task_id, "role": role, "title": title, "status": "error",
                "summary": msg, "tools": []}
    except Exception as e:
        msg = f"{type(e).__name__}: {e}"
        try:
            await emit_queue.put({"type": "agent_error", "taskId": task_id, "error": msg})
        except Exception:
            pass
        return {"taskId": task_id, "role": role, "title": title, "status": "error",
                "summary": msg, "tools": []}


async def _orchestrator_run_agent_inner(
    gw, stream_model: str, subtask: dict, project_path: str,
    base_system_prompt: str, emit_queue: asyncio.Queue,
    aws_profile: str = "", bedrock_user: str = "", is_remote: bool = False,
    template_id: str = "",
):
    """실제 에이전트 실행 본체."""
    task_id = subtask.get("id", "?")
    role = subtask.get("role", "Worker")
    title = subtask.get("title", "")
    description = subtask.get("description", "")
    target_files = subtask.get("target_files", [])
    primary_tool = subtask.get("primary_tool", "")

    sys_prompt = ORCHESTRATOR_AGENT_PROMPT.format(
        role=role, task_id=task_id, title=title,
        target_files=", ".join(target_files) if target_files else "(없음)",
        primary_tool=primary_tool or "(자동 선택)",
        description=description,
    )
    if base_system_prompt:
        sys_prompt = base_system_prompt + "\n\n" + sys_prompt

    messages = [{"role": "user", "content": [{"text": f"작업 {task_id} ({role}): {title}\n\n{description}"}]}]
    max_turns = int(os.environ.get("AE_MAX_ORCH_TURNS", "50"))
    # 같은 도구가 무한 반복 호출되는 것을 막기 위한 cap.
    # 사용자 환경에서 14m 58s 동안 0건 생성되며 도구 호출 폭주가 발견됨 (13/9/5/3/3회).
    # AE_MAX_TOOL_CALLS_PER_AGENT 환경변수로 조정 가능, 기본 8회.
    max_tool_calls_per_agent = int(os.environ.get("AE_MAX_TOOL_CALLS_PER_AGENT", "8"))
    tool_call_count = 0
    final_text_parts = []
    tool_log = []
    # 실제 디스크에 존재 확인된 파일들 — 할루시네이션 방지의 핵심
    verified_files = []
    # Worker가 도구 안 부르고 텍스트만 답할 때 자동 재요청 카운터 (max 2회)
    tool_nudge_count = 0
    MAX_TOOL_NUDGES = 2

    await emit_queue.put({"type": "agent_start", "taskId": task_id, "role": role, "title": title, "targetFiles": target_files})

    # === 결정론적 단축 경로 (NEW) ===
    # primary_tool이 명확한 파일 생성 도구면 Worker LLM 호출 없이 바로 enrich+generate.
    # 이러면 Worker가 도구를 안 부르는 모든 케이스를 우회한다.
    # 게이트웨이는 enrichment 단계에서 1회 경유 → 비용/사용량 측정 유지.
    pt_lower = (subtask.get("primary_tool") or "").lower()
    direct_tools = {"generate_pdf", "generate_pptx", "generate_xlsx",
                    "generate_docx", "generate_image", "edit_image", "write_file"}
    if pt_lower in direct_tools:
        print(f"[Orchestrator] {task_id} 결정론적 단축 경로 — Worker LLM 우회 (primary_tool={pt_lower})")
        try:
            # 1) 게이트웨이로 Claude가 콘텐츠 작성 (게이트웨이 경유 보장).
            #    enrich task → Sonnet 라우팅 (stream_model이 사용자 적합 모델이면 그대로).
            _enrich_model = _specialized_model_for_task(
                "enrich", stream_model,
                aws_profile=aws_profile, bedrock_user=bedrock_user,
            )
            enriched_text = await _enrich_content_via_gateway(
                gw=gw,
                model_id=_enrich_model,
                primary_tool=pt_lower,
                title=title,
                description=description,
                final_text="",
                project_path=project_path,
            )
            # 2) enrichment 부족하면 디스크 데이터 보강
            if not enriched_text or len(enriched_text) < 200:
                real = _gather_real_context(description, project_path) \
                    or _gather_real_context_forced(project_path, title)
                if real:
                    enriched_text = real
            # 3) 결정론적 도구로 파일 생성
            forced = await _force_generate_from_text(
                primary_tool=pt_lower,
                target_files=target_files,
                title=title,
                description=description,
                final_text=enriched_text or description or title,
                project_path=project_path,
                aws_profile=aws_profile,
                bedrock_user=bedrock_user,
                template_id=template_id,
            )
            for fpath, finfo in forced:
                verified_files.append(finfo)
                tool_log.append({
                    "name": "direct:" + finfo.get("tool", "?"),
                    "input": {"shortcut": True, "primary_tool": pt_lower},
                    "output": fpath,
                })
                # meta sidecar
                try:
                    _meta = {
                        "tool": finfo.get("tool", pt_lower),
                        "model": finfo.get("model", "system_direct"),
                        "chatModel": stream_model,
                        "agentId": task_id,
                        "agentRole": role,
                        "agentTitle": title,
                        "createdAt": datetime.utcnow().isoformat() + "Z",
                        "promptHint": (description or "")[:200],
                        "shortcut": True,
                    }
                    with open(finfo["absPath"] + ".meta.json", "w", encoding="utf-8") as _mf:
                        json.dump(_meta, _mf, ensure_ascii=False, indent=2)
                except Exception:
                    pass
            if verified_files:
                final_text = (enriched_text or "")[:1500]
                await emit_queue.put({
                    "type": "agent_done", "taskId": task_id,
                    "summary": f"단축 경로로 {len(verified_files)}개 파일 생성됨",
                    "toolCount": len(tool_log),
                    "verifiedFiles": [vf["path"] for vf in verified_files],
                })
                return {
                    "taskId": task_id, "role": role, "title": title,
                    "status": "done",
                    "summary": f"[결정론적 생성] {pt_lower} 도구로 파일 생성 완료\n\n{final_text}",
                    "tools": tool_log, "verifiedFiles": verified_files,
                }
            print(f"[Orchestrator] {task_id} 단축 경로 실패 — Worker LLM 경로로 진행")
        except Exception as _short_err:
            print(f"[Orchestrator] {task_id} 단축 경로 예외: {_short_err}, Worker LLM으로 진행")

    stream_failed = False
    stream_error_msg = ""

    try:
        for turn in range(max_turns):
            text_parts = []
            tool_use_blocks = []
            current_tool = {}
            stop_reason = ""

            try:
                async for evt in gw.stream_sse_realtime(
                    model_id=stream_model, messages=messages,
                    system_prompt=sys_prompt, tool_config=AGENT_TOOLS,
                ):
                    etype = evt.get("type", "")
                    if etype == "content_block_delta":
                        delta = evt.get("delta", {})
                        if "text" in delta:
                            text_parts.append(delta["text"])
                            await emit_queue.put({"type": "agent_delta", "taskId": task_id, "text": delta["text"]})
                        elif "toolUse" in delta:
                            if current_tool:
                                current_tool["_input_json"] = current_tool.get("_input_json", "") + delta["toolUse"].get("input", "")
                    elif etype == "content_block_start":
                        cb = evt.get("content_block") or evt.get("contentBlock") or {}
                        if "toolUse" in cb:
                            tu = cb["toolUse"]
                            current_tool = {"toolUseId": tu.get("toolUseId", ""), "name": tu.get("name", ""), "_input_json": ""}
                    elif etype == "content_block_stop":
                        if current_tool and current_tool.get("name"):
                            try:
                                inp = json.loads(current_tool.get("_input_json", "{}"))
                            except json.JSONDecodeError:
                                inp = {}
                            tool_use_blocks.append({"toolUse": {"toolUseId": current_tool["toolUseId"], "name": current_tool["name"], "input": inp}})
                            current_tool = {}
                    elif etype in ("message_delta", "message_stop"):
                        stop_reason = evt.get("delta", {}).get("stopReason", "") or evt.get("stop_reason", "") or evt.get("stopReason", "")
                    elif etype == "error":
                        await emit_queue.put({"type": "agent_error", "taskId": task_id, "error": evt.get("message", "")})
            except Exception as e:
                stream_failed = True
                stream_error_msg = str(e)
                await emit_queue.put({"type": "agent_error", "taskId": task_id, "error": stream_error_msg})
                break

            content_blocks = []
            if text_parts:
                content_blocks.append({"text": "".join(text_parts)})
                final_text_parts.extend(text_parts)
            content_blocks.extend(tool_use_blocks)
            if not content_blocks:
                break
            messages.append({"role": "assistant", "content": content_blocks})
            if not tool_use_blocks:
                if stop_reason == "max_tokens" and text_parts and turn < max_turns - 1:
                    print(f"[Orchestrator] max_tokens 도달 — 이어서 생성 (task={task_id}, turn={turn+1})")
                    messages.append({"role": "user", "content": [{"text": "계속 이어서 작성해주세요."}]})
                    continue
                # === 약점 1 개선: 도구 호출 0회 자동 nudge ===
                # 작업이 도구를 필요로 하는데 Worker가 텍스트만 답함 → 도구 호출 재요청.
                # forced fallback에 의존하지 않고 정상 경로로 성공 확률 향상.
                pt = (subtask.get("primary_tool") or "").lower()
                target_files = subtask.get("target_files", [])
                wants_tools = bool(target_files) or pt in (
                    "generate_image", "generate_pdf", "generate_pptx",
                    "generate_xlsx", "generate_docx", "edit_image",
                    "write_file", "code", "run_command",
                )
                if wants_tools and tool_nudge_count < MAX_TOOL_NUDGES and turn < max_turns - 1:
                    tool_nudge_count += 1
                    nudge_msg = (
                        f"위 응답에서 도구 호출을 하지 않았습니다.\n"
                        f"이 작업은 반드시 도구로 실제 파일을 만들어야 합니다.\n"
                        f"primary_tool={pt or '(자동선택)'} 도구를 즉시 호출해서 "
                        f".generated/ 폴더에 파일을 생성해주세요.\n"
                        f"텍스트 응답만으로는 작업이 실패한 것으로 간주됩니다."
                    )
                    if pt == "generate_pdf":
                        nudge_msg += '\n예시: generate_pdf({"title": "...", "sections": [{"heading": "...", "body": "..."}]})'
                    elif pt == "generate_pptx":
                        nudge_msg += '\n예시: generate_pptx({"title": "...", "slides": [{"title": "...", "bullets": ["..."]}]})'
                    elif pt == "generate_xlsx":
                        nudge_msg += '\n예시: generate_xlsx({"title": "...", "sheets": [{"name": "Data", "headers": ["A", "B"], "rows": [["x", "y"]]}]})'
                    elif pt == "generate_docx":
                        nudge_msg += '\n예시: generate_docx({"title": "...", "sections": [{"heading": "...", "body": "...", "bullets": ["..."]}]})'
                    elif pt == "generate_image":
                        nudge_msg += '\n예시: generate_image({"prompt": "...", "size": "1024x1024"})'
                    messages.append({"role": "user", "content": [{"text": nudge_msg}]})
                    print(f"[Orchestrator] {task_id} 도구 nudge {tool_nudge_count}/{MAX_TOOL_NUDGES} (primary_tool={pt})")
                    await emit_queue.put({
                        "type": "agent_tool", "taskId": task_id,
                        "tool": "tool_nudge", "status": "running",
                        "input": {"attempt": tool_nudge_count, "primary_tool": pt},
                    })
                    continue
                break

            tool_results = []
            for block in tool_use_blocks:
                tu = block["toolUse"]
                tname = tu.get("name", "")
                tid = tu.get("toolUseId", "")
                tinput = tu.get("input", {})
                # === 도구 호출 cap ===
                # 모델이 같은 도구를 무한 재호출해 14분+ 폭주하는 것을 방지.
                # cap 도달 시 즉시 단락하고 toolResult로 cap 알림 → 모델이 텍스트로
                # 답변을 마무리하도록 유도. 그래도 verified_files=0이면 외부 forced
                # fallback이 디스크에 직접 파일을 생성한다.
                tool_call_count += 1
                if tool_call_count > max_tool_calls_per_agent:
                    cap_msg = (
                        f"[도구 호출 한도 초과] 이 에이전트는 {max_tool_calls_per_agent}회를 초과해 "
                        f"호출했습니다 (현재 {tool_call_count}회). 더 이상 도구를 호출하지 말고 "
                        f"지금까지 한 작업을 한 줄로 요약해 답변을 마무리하세요."
                    )
                    print(f"[Orchestrator] {task_id} 도구 호출 cap 도달 ({tool_call_count}/{max_tool_calls_per_agent}) — {tname} 단락")
                    await emit_queue.put({
                        "type": "agent_tool", "taskId": task_id,
                        "tool": tname, "status": "capped",
                        "input": {"reason": "tool_call_cap", "limit": max_tool_calls_per_agent},
                    })
                    tool_log.append({"name": tname, "input": tinput, "output": cap_msg, "capped": True})
                    tool_results.append({"toolResult": {"toolUseId": tid, "content": [{"text": cap_msg}]}})
                    continue
                await emit_queue.put({"type": "agent_tool", "taskId": task_id, "tool": tname, "input": tinput, "status": "running"})
                tout = await asyncio.to_thread(_execute_tool, tname, tinput, project_path, aws_profile, bedrock_user, template_id)  # [patched-credentials]
                tool_log.append({"name": tname, "input": tinput, "output": tout[:400]})

                # 생성 파일에 .meta.json 사이드카 기록 — 어떤 모델/에이전트가 만들었는지 추적
                if tname in ("generate_image", "generate_pdf", "generate_pptx", "generate_xlsx", "generate_docx", "edit_image", "write_file"):
                    try:
                        _meta_paths = []
                        _actual_model = ""  # 실제 생성에 사용된 모델 (도구 응답에서 추출)
                        try:
                            _parsed = json.loads(tout) if isinstance(tout, str) else None
                        except (json.JSONDecodeError, TypeError):
                            _parsed = None
                        if isinstance(_parsed, dict) and "path" in _parsed and "error" not in _parsed:
                            _meta_paths.append(_parsed["path"])
                            _actual_model = _parsed.get("model") or ""
                        elif isinstance(_parsed, dict) and isinstance(_parsed.get("images"), list):
                            for _it in _parsed["images"]:
                                if isinstance(_it, dict) and "path" in _it:
                                    _meta_paths.append(_it["path"])
                                    if not _actual_model and _it.get("model"):
                                        _actual_model = _it["model"]
                        # write_file: tool_input["path"]
                        if tname == "write_file" and "path" in tinput:
                            _meta_paths.append(tinput["path"])
                        # 도구 카테고리별 기본 모델명 (도구 응답에 model이 없을 때)
                        _tool_default_model = {
                            "generate_pdf":  "reportlab (Python)",
                            "generate_pptx": "python-pptx",
                            "generate_xlsx": "openpyxl",
                            "generate_docx": "python-docx",
                            "write_file":    "filesystem",
                        }
                        for _rel in _meta_paths:
                            # 도구가 _resolve_local_root로 저장한 위치를 그대로 따라감.
                            # 단순 project_path join은 AE_GENERATED_ROOT/~/.agentic-editor
                            # fallback에 저장된 파일을 못 찾아 verified_files=0 버그 유발했음.
                            _abs = _resolve_relative_for_verify(_rel, project_path)
                            if not os.path.isfile(_abs):
                                continue
                            # 실제 디스크 검증 통과 → verified_files에 추가
                            try:
                                _stat = os.stat(_abs)
                                verified_files.append({
                                    "path": _rel,
                                    "absPath": _abs,
                                    "size": _stat.st_size,
                                    "tool": tname,
                                    "model": _actual_model or "",
                                })
                            except Exception:
                                pass
                            # 실제 사용 모델 = (이미지 생성 모델) 또는 (도구별 기본) 또는 (chat 모델)
                            _model_label = _actual_model or _tool_default_model.get(tname) or stream_model
                            _meta_obj = {
                                "tool": tname,
                                "model": _model_label,           # 실제 작업 수행 모델/엔진
                                "chatModel": stream_model,       # 도구 호출을 결정한 chat 모델
                                "agentId": task_id,
                                "agentRole": role,
                                "agentTitle": title,
                                "createdAt": datetime.utcnow().isoformat() + "Z",
                                "promptHint": (subtask.get("description") or "")[:200],
                            }
                            try:
                                with open(_abs + ".meta.json", "w", encoding="utf-8") as _mf:
                                    json.dump(_meta_obj, _mf, ensure_ascii=False, indent=2)
                            except Exception as _me:
                                print(f"[Orchestrator] meta write 실패 {_abs}: {_me}")
                    except Exception as _e:
                        print(f"[Orchestrator] meta 처리 예외: {_e}")

                await emit_queue.put({"type": "agent_tool", "taskId": task_id, "tool": tname, "status": "done", "output": tout[:300]})
                _tr_max = int(os.environ.get("AE_TOOL_RESULT_MAX", "80000"))
                tool_results.append({"toolResult": {"toolUseId": tid, "content": [{"text": tout[:_tr_max]}]}})
            messages.append({"role": "user", "content": tool_results})

        final_text = "".join(final_text_parts).strip()
        if stream_failed:
            # SSE 실패해도 forced fallback은 시도 — agent_error가 이미 emit되었지만
            # 사용자가 파일을 요구했으면 디스크 데이터로라도 만들어준다.
            primary_tool = (subtask.get("primary_tool") or "").lower()
            wanted_files = bool(target_files) or primary_tool in (
                "generate_image", "generate_pdf", "generate_pptx",
                "generate_xlsx", "generate_docx", "edit_image", "write_file"
            )
            if wanted_files:
                print(f"[Orchestrator] {task_id} SSE 실패 → forced fallback 시도")
                try:
                    real_only = _gather_real_context(description, project_path) \
                        or _gather_real_context_forced(project_path, title)
                    if real_only:
                        _forced = await _force_generate_from_text(
                            primary_tool=primary_tool,
                            target_files=target_files,
                            title=title,
                            description=description,
                            final_text=real_only,
                            project_path=project_path,
                            aws_profile=aws_profile,
                            bedrock_user=bedrock_user,
                            template_id=template_id,
                        )
                        for fpath, finfo in _forced:
                            verified_files.append(finfo)
                            tool_log.append({
                                "name": finfo.get("tool", "deterministic"),
                                "input": {"forced": True, "reason": "stream_failed"},
                                "output": fpath,
                            })
                        if verified_files:
                            await emit_queue.put({
                                "type": "agent_done", "taskId": task_id,
                                "summary": f"SSE 실패 후 강제 생성 완료 ({len(verified_files)}개 파일)",
                                "toolCount": len(tool_log),
                                "verifiedFiles": [vf["path"] for vf in verified_files],
                            })
                            return {
                                "taskId": task_id, "role": role, "title": title,
                                "status": "done",
                                "summary": f"[강제 생성] {stream_error_msg or '스트리밍 실패'} → 디스크 데이터로 파일 생성됨",
                                "tools": tool_log, "verifiedFiles": verified_files,
                            }
                except Exception as _e:
                    print(f"[Orchestrator] SSE 실패 후 강제 생성도 실패: {_e}")
            return {"taskId": task_id, "role": role, "title": title, "status": "error",
                    "summary": stream_error_msg or "스트리밍 실패", "tools": tool_log,
                    "verifiedFiles": []}

        # === Worker bypass detection (Change 3) ===
        # primary_tool이 generate_pdf 인데 worker가 write_file로 .py/.html/.md 만들었거나,
        # primary_tool이 generate_pptx 인데 generate_image로 .png/.svg 만든 경우 등을
        # 잡아낸다. 비매칭 파일은 verified_files에 그대로 두되 primary 형식 요구는
        # 충족되지 않은 것으로 간주 → 아래 wanted_files 분기에서 forced fallback이
        # 정확히 required_ext로 강제 생성되도록 한다.
        bypass_detected = False
        _required_ext = {
            "generate_pdf": "pdf",
            "generate_pptx": "pptx",
            "generate_xlsx": "xlsx",
            "generate_docx": "docx",
            "generate_image": "png",
        }.get((subtask.get("primary_tool") or "").lower(), "")
        if _required_ext and verified_files:
            _matching = [vf for vf in verified_files
                         if str(vf.get("path", "")).lower().endswith(f".{_required_ext}")]
            _bypass = [vf for vf in verified_files
                       if not str(vf.get("path", "")).lower().endswith(f".{_required_ext}")]
            if not _matching and _bypass:
                print(
                    f"[Orchestrator] {task_id} 우회 감지 — "
                    f"primary_tool={subtask.get('primary_tool')}, "
                    f"생성된 비매칭 파일: {[vf['path'] for vf in _bypass]}"
                )
                bypass_detected = True
                # bypass 파일은 verified_files에 남겨두지만, 아래 wanted_files 분기가
                # _matching 0건임을 보고 forced fallback을 트리거할 수 있도록
                # verified_files 자체는 _bypass 그대로 (=0 매칭).
                verified_files = _bypass
                try:
                    await emit_queue.put({
                        "type": "agent_tool", "taskId": task_id,
                        "tool": "bypass_detector", "status": "done",
                        "input": {"primary_tool": subtask.get("primary_tool"),
                                  "expected_ext": _required_ext},
                        "output": f"비매칭 파일 {len(_bypass)}건, 매칭 0건 → 강제 재생성 예정",
                    })
                except Exception:
                    pass

        # 파일 생성 작업이었는지 판정 — primary_tool이나 target_files로 판단
        primary_tool = (subtask.get("primary_tool") or "").lower()
        wanted_files = bool(target_files) or primary_tool in (
            "generate_image", "generate_pdf", "generate_pptx",
            "generate_xlsx", "generate_docx", "edit_image", "write_file"
        )
        # description에 파일 형식 키워드가 있으면 wanted_files=True
        # (Directory Analyzer 같은 generic role도 fallback 진입 가능)
        desc_lower = (description or "").lower() + " " + (title or "").lower()
        if not wanted_files:
            # 진짜 파일 생성을 의미하는 강한 키워드만 — "분석", "리뷰" 같은 일반 단어는 제외
            file_keywords = ("pdf", "pptx", "xlsx", "docx", "png", "jpg",
                             "엑셀", "워드", "파워포인트", "프레젠테이션",
                             "스프레드시트", "보고서", "report.pdf", "report.docx")
            # 단, "분석/리뷰/검토" 같은 일반 단어가 단독으로 있으면 forced 발동 안 함
            soft_keywords = ("분석", "리뷰", "검토", "구조", "흐름도", "다이어그램", "차트", "그래프", "이미지", "문서")
            has_strong = any(kw in desc_lower for kw in file_keywords)
            has_soft = any(kw in desc_lower for kw in soft_keywords)
            # 강한 키워드만 forced 진입 — soft만 있으면 모델이 결정한 결과 신뢰 (텍스트 응답 가능)
            if has_strong:
                wanted_files = True
                # primary_tool 추론 — title/description에서 형식 감지
                if not primary_tool:
                    if "pdf" in desc_lower:
                        primary_tool = "generate_pdf"
                    elif "pptx" in desc_lower or "프레젠테이션" in desc_lower or "파워포인트" in desc_lower:
                        primary_tool = "generate_pptx"
                    elif "xlsx" in desc_lower or "엑셀" in desc_lower or "스프레드시트" in desc_lower:
                        primary_tool = "generate_xlsx"
                    elif "docx" in desc_lower or "워드" in desc_lower:
                        primary_tool = "generate_docx"
                    elif "png" in desc_lower or "jpg" in desc_lower:
                        primary_tool = "generate_image"
                    else:
                        primary_tool = "generate_pdf"
                    print(f"[Orchestrator] {task_id} primary_tool 자동 감지 → {primary_tool}")
            elif has_soft:
                # soft 키워드만 — 텍스트 응답으로도 충분, forced 스킵
                print(f"[Orchestrator] {task_id} soft keyword only → forced 스킵 (텍스트 응답 신뢰)")

        # 할루시네이션 차단 + 강제 생성:
        # 파일 생성을 요구받았는데 실제 디스크에 검증된 파일이 0개면 →
        # (1) 게이트웨이 1회 추가 호출로 본문 보강 → (2) 시스템이 도구 디스패처로 파일 강제 생성.
        # 게이트웨이 경유 유지 (비용/사용량 측정).
        if wanted_files and not verified_files:
            print(f"[Orchestrator] {task_id} 도구 미호출 감지 — 게이트웨이 보강 + 강제 fallback (primary_tool={primary_tool})")
            # bypass_detected면 target_files 확장자를 required_ext로 강제 교체.
            # 이 경로는 worker가 .html/.py 등으로 우회한 케이스를 정상 형식으로 복구한다.
            forced_target_files = list(target_files or [])
            if bypass_detected and _required_ext:
                _slug = re.sub(r"[^A-Za-z0-9._-]+", "-",
                               (title or task_id or "result").strip())[:40] or "result"
                forced_target_files = [f".generated/forced-{_slug}.{_required_ext}"]
                print(f"[Orchestrator] {task_id} bypass 복구 — 강제 target_files={forced_target_files}")
            await emit_queue.put({
                "type": "agent_tool", "taskId": task_id,
                "tool": "deterministic-converter", "status": "running",
                "input": {"reason": "no tool calls — gateway enrich + force generate"},
            })
            try:
                # (1) 게이트웨이로 콘텐츠 보강 — 짧은 final_text를 풍부한 본문으로 확장.
                #     stream_model은 이미 자동 라우팅되어 도구 호출 가능 모델.
                #     enrich task → Sonnet 우선 라우팅.
                _enrich_model = _specialized_model_for_task(
                    "enrich", stream_model,
                    aws_profile=aws_profile, bedrock_user=bedrock_user,
                )
                enriched_text = await _enrich_content_via_gateway(
                    gw=gw,
                    model_id=_enrich_model,
                    primary_tool=primary_tool,
                    title=title,
                    description=description,
                    final_text=final_text,
                    project_path=project_path,
                )

                # 게이트웨이 enrichment가 비어있거나 너무 짧으면 실제 디스크 데이터를
                # 직접 final_text로 사용 — Claude 호출 없이도 파일은 무조건 생성됨.
                if not enriched_text or len(enriched_text) < 200:
                    real_only = _gather_real_context(description, project_path)
                    if not real_only:
                        # 키워드 매칭 안 됨 → forced 모드로 폴더+코드 인벤토리 수집
                        real_only = _gather_real_context_forced(project_path, title)
                    if real_only:
                        enriched_text = real_only
                        print(f"[Orchestrator] {task_id} enrichment 부족 — 실제 디스크 데이터로 대체 ({len(real_only)}자)")
                # (2) 보강된 본문으로 결정적 도구 디스패처 호출 → 실제 파일 생성
                _forced = await _force_generate_from_text(
                    primary_tool=primary_tool,
                    target_files=forced_target_files,
                    title=title,
                    description=description,
                    final_text=enriched_text,
                    project_path=project_path,
                    aws_profile=aws_profile,
                    bedrock_user=bedrock_user,
                    template_id=template_id,
                )
                for fpath, finfo in _forced:
                    verified_files.append(finfo)
                    tool_log.append({
                        "name": finfo.get("tool", "deterministic"),
                        "input": {"forced": True},
                        "output": fpath,
                    })
                    # meta sidecar
                    try:
                        _meta_obj = {
                            "tool": finfo.get("tool", "deterministic"),
                            # forced 경로에서도 실제 사용 라이브러리 이름이 박힘
                            # (reportlab/python-pptx/python-docx/openpyxl/matplotlib).
                            # system_fallback 라벨은 절대 사용 안 함.
                            "model": finfo.get("model", "deterministic-converter"),
                            "chatModel": stream_model,
                            "agentId": task_id,
                            "agentRole": role,
                            "agentTitle": title,
                            "createdAt": datetime.utcnow().isoformat() + "Z",
                            "promptHint": "[forced fallback] " + (description or "")[:180],
                            "forced": True,
                        }
                        with open(finfo["absPath"] + ".meta.json", "w", encoding="utf-8") as _mf:
                            json.dump(_meta_obj, _mf, ensure_ascii=False, indent=2)
                    except Exception:
                        pass
                await emit_queue.put({
                    "type": "agent_tool", "taskId": task_id,
                    "tool": "deterministic-converter", "status": "done",
                    "output": f"forced {len(_forced)} files",
                })
            except Exception as _fe:
                err_msg = f"강제 생성 실패: {_fe}"
                print(f"[Orchestrator] {task_id} {err_msg}")
                await emit_queue.put({
                    "type": "agent_error", "taskId": task_id,
                    "error": err_msg,
                })
                return {
                    "taskId": task_id, "role": role, "title": title,
                    "status": "error",
                    "summary": f"{final_text}\n\n[시스템] {err_msg}",
                    "tools": tool_log, "verifiedFiles": [],
                }

        # 강제 fallback 후에도 파일이 0개면 정말 실패
        if wanted_files and not verified_files:
            err_msg = "파일 생성 실패 — 강제 fallback도 결과 0건"
            print(f"[Orchestrator] {task_id} {err_msg}")
            await emit_queue.put({
                "type": "agent_error", "taskId": task_id,
                "error": err_msg,
            })
            return {
                "taskId": task_id, "role": role, "title": title,
                "status": "error",
                "summary": f"{final_text}\n\n[시스템] {err_msg}",
                "tools": tool_log, "verifiedFiles": [],
            }

        await emit_queue.put({
            "type": "agent_done", "taskId": task_id,
            "summary": final_text[-1200:],
            "toolCount": len(tool_log),
            "verifiedFiles": [vf["path"] for vf in verified_files],
        })
        return {
            "taskId": task_id, "role": role, "title": title,
            "status": "done", "summary": final_text, "tools": tool_log,
            "verifiedFiles": verified_files,
        }
    except Exception as e:
        await emit_queue.put({"type": "agent_error", "taskId": task_id, "error": str(e)})
        return {"taskId": task_id, "role": role, "title": title, "status": "error",
                "summary": str(e), "tools": tool_log, "verifiedFiles": []}


def _deterministic_failure_report(user_prompt: str, agent_results: list) -> str:
    """모든 에이전트가 도구 호출 0회로 실패했을 때의 정직한 보고서.

    LLM 없이 직접 생성 → KeyError 같은 가짜 오류 등장 불가.
    """
    rows = []
    for r in agent_results:
        tid = r.get("taskId", "?")
        role = r.get("role", "Worker")
        rows.append(f"| {tid} | {role} | 실패 | 없음 | 0 |")
    table = "\n".join(rows)
    return f"""## 최종 통합 결과

| 에이전트 | 역할 | 상태 | 생성된 파일 (디스크 검증됨) | 도구 사용 횟수 |
|---------|------|------|--------------------------|--------------|
{table}

### 생성된 파일 목록 (verifiedFiles 기반 — 실제 존재)

⚠️ 디스크에 검증된 파일이 0개입니다. 어떤 파일도 실제로 생성되지 않았습니다.

### 세부 사항

모든 에이전트가 도구를 한 번도 호출하지 않았습니다. 다음 원인 중 하나가 가능성 높습니다:

- 워커 모델이 toolConfig를 무시하고 텍스트로만 응답함
- 게이트웨이가 toolConfig를 워커 모델로 전달하지 못함
- 강제 fallback 경로에서 디스크 데이터 수집이 실패함

### 권장 조치

1. **서버 재시작** — 코드 변경이 반영되지 않았을 수 있습니다.
2. **다른 모델로 재시도** — Claude Sonnet 4.6 / Opus 4.7 권장.
3. **`.generated/` 폴더 확인** — 강제 fallback이 부분 성공했을 수 있습니다.
"""


def _deterministic_success_report(user_prompt: str, agent_results: list, verified_files: list,
                                  required_formats: set = None) -> str:
    """모든 에이전트 성공 + 검증된 파일 있을 때의 결정론적 보고서.

    required_formats가 주어지면 사용자 명시 형식 충족도를 명시한다 — 누락이
    있으면 보고서 끝에 "**누락**: ..." 섹션이 추가되어 false-positive를 차단.
    """
    required_formats = required_formats or set()
    agent_rows = []
    for r in agent_results:
        tid = r.get("taskId", "?")
        role = r.get("role", "Worker")
        files_for_agent = [vf for vf in verified_files if vf.get("agentId") == tid]
        files_str = ", ".join(f"`{vf['path']}`" for vf in files_for_agent) or "없음"
        tool_count = len(r.get("tools", []))
        agent_rows.append(f"| {tid} | {role} | ✓ 완료 | {files_str} | {tool_count} |")

    file_lines = []
    for vf in verified_files:
        path = vf.get("path", "?")
        size = vf.get("size", 0)
        model = vf.get("model", "")
        agent = vf.get("agentRole", "")
        size_kb = f"{size / 1024:.1f} KB" if size >= 1024 else f"{size} B"
        meta = f" — {agent}" if agent else ""
        if model:
            meta += f" ({model})"
        file_lines.append(f"- `{path}` ({size_kb}){meta}")

    # 사용자 요청 형식 충족도 계산 (요청이 있을 때만)
    coverage_section = ""
    if required_formats:
        have = set()
        for vf in verified_files:
            p = (vf.get("path") or "").lower()
            for ext in ("pdf", "pptx", "xlsx", "docx", "png", "jpg", "jpeg"):
                if p.endswith(f".{ext}"):
                    have.add("png" if ext in ("jpg", "jpeg") else ext)
                    break
        missing = required_formats - have
        if missing:
            coverage_section = (
                f"\n\n### 사용자 요청 형식 충족도\n"
                f"- 요청: {sorted(required_formats)}\n"
                f"- 생성됨: {sorted(have & required_formats)}\n"
                f"- **누락**: {sorted(missing)}\n"
            )
        else:
            coverage_section = (
                f"\n\n### 사용자 요청 형식 충족도\n"
                f"- 모든 요청 형식 생성됨: {sorted(required_formats)}\n"
            )

    return f"""## 최종 통합 결과

| 에이전트 | 역할 | 상태 | 생성된 파일 (디스크 검증됨) | 도구 사용 횟수 |
|---------|------|------|--------------------------|--------------|
{chr(10).join(agent_rows)}

### 생성된 파일 목록 (verifiedFiles 기반 — 실제 존재)

{chr(10).join(file_lines)}

### 세부 사항

총 {len(agent_results)}개 에이전트가 모두 작업을 완료했고, 디스크에 {len(verified_files)}개 파일이 검증되었습니다.{coverage_section}
"""


async def _orchestrator_merge(gw, stream_model, user_prompt, agent_results: list, base_system_prompt: str,
                              required_formats: set = None) -> str:
    """Merger 호출 — 최종 보고서 생성. verifiedFiles 기반으로 거짓 완료 주장 차단.

    required_formats: 사용자가 명시한 파일 형식들. 결정론적 성공 보고서와
    LLM 기반 Merger 모두에 전달되어 누락 형식이 솔직히 보고되도록 한다.
    """
    required_formats = required_formats or set()
    # 실제로 디스크에 검증된 파일들만 추출
    all_verified_files = []
    for r in agent_results:
        for vf in r.get("verifiedFiles", []) or []:
            if isinstance(vf, dict) and vf.get("path"):
                all_verified_files.append({
                    "agentId": r.get("taskId"),
                    "agentRole": r.get("role"),
                    "path": vf["path"],
                    "size": vf.get("size", 0),
                    "tool": vf.get("tool", "?"),
                    "model": vf.get("model", ""),
                })

    # === 결정론적 short-circuit ===
    # 모든 에이전트가 실패 + 도구 호출 0회면 LLM 호출 없이 정직한 보고서 직접 생성.
    # LLM이 가짜 KeyError 등을 만들어내는 것을 원천 차단.
    all_failed = all(r.get("status") != "done" for r in agent_results) if agent_results else True
    no_tools = all(len(r.get("tools", [])) == 0 for r in agent_results) if agent_results else True
    no_files = len(all_verified_files) == 0
    if all_failed and no_tools and no_files and agent_results:
        return _deterministic_failure_report(user_prompt, agent_results)

    # === 모든 에이전트 성공 + 검증된 파일 있음 ===
    # LLM 없이 결정론적 성공 보고서 — Merger 할루시네이션 위험 제거.
    all_done = all(r.get("status") == "done" for r in agent_results) if agent_results else False
    if all_done and all_verified_files and agent_results:
        return _deterministic_success_report(user_prompt, agent_results, all_verified_files,
                                             required_formats=required_formats)

    summary_input = {
        "userRequest": user_prompt[:2000],
        "requiredFormats": sorted(required_formats) if required_formats else [],
        "agents": [
            {
                "taskId": r["taskId"], "role": r["role"], "title": r["title"],
                "status": r["status"],
                "summary": (r.get("summary") or "")[:1500],
                "toolCount": len(r.get("tools", [])),
                "verifiedFileCount": len(r.get("verifiedFiles") or []),
                # 실제 도구 호출 명단 — Merger가 가짜 오류 만들어내지 않도록
                "toolsCalled": [t.get("name", "?") for t in (r.get("tools") or [])][:10],
                # 도구 호출 0회면 분명히 표시 — Merger가 추측 금지
                "noToolsCalled": (len(r.get("tools", [])) == 0),
            }
            for r in agent_results
        ],
        "verifiedFiles": all_verified_files,  # 실제 디스크 검증된 파일만
    }
    sys_prompt = ORCHESTRATOR_MERGER_PROMPT
    if base_system_prompt:
        sys_prompt = base_system_prompt + "\n\n" + sys_prompt
    messages = [{"role": "user", "content": [{"text": json.dumps(summary_input, ensure_ascii=False, indent=2)}]}]
    try:
        result = await asyncio.wait_for(
            gw.converse(model_id=stream_model, messages=messages, system_prompt=sys_prompt),
            timeout=120,
        )
        if result.get("decision") == "ALLOW":
            output = result.get("output", {}).get("message", {}).get("content", [])
            text = "\n".join(c.get("text", "") for c in output if "text" in c).strip()
            # === 약점 3 개선: Merger 응답 sanitize ===
            # 입력에 없는 가짜 오류명(KeyError, TypeError 등)을 만들어내면 제거.
            return _sanitize_merger_report(text, agent_results)
        return f"(Merger 실패: {result.get('error') or result.get('decision')})"
    except Exception as e:
        return f"(Merger 예외: {e})"


def _sanitize_merger_report(text: str, agent_results: list) -> str:
    """Merger 응답에서 입력에 없는 가짜 예외명/스택트레이스 제거.

    LLM이 'KeyError', 'TypeError', 'ValidationException' 같은 그럴듯한 오류 메시지를
    만들어내는 경우가 있음. 실제 입력(agent.summary)에 그 단어가 없으면 일반화된
    문구로 치환한다.
    """
    if not text:
        return text

    # 실제 agent summary에 등장한 예외명들 수집 (이건 OK — 진짜 정보)
    real_errors = set()
    for r in (agent_results or []):
        s = (r.get("summary") or "").lower()
        for kw in ("keyerror", "typeerror", "valueerror", "attributeerror",
                   "validationexception", "accessdeniedexception",
                   "filenotfound", "permissiondenied"):
            if kw in s:
                real_errors.add(kw)

    # 입력에 없는 fake 예외명 패턴
    fake_patterns = [
        (r"KeyError:\s*['\"][^'\"]*['\"]", "keyerror"),
        (r"TypeError:\s*['\"][^'\"]*['\"]", "typeerror"),
        (r"ValueError:\s*['\"][^'\"]*['\"]", "valueerror"),
        (r"ValidationException[^.\n]*", "validationexception"),
        (r"AttributeError:\s*['\"][^'\"]*['\"]", "attributeerror"),
    ]
    sanitized = text
    for pattern, kw in fake_patterns:
        if kw not in real_errors:
            sanitized = re.sub(pattern, "도구 호출 없이 텍스트만 출력함", sanitized)

    return sanitized


@app.post("/api/agents/run-orchestrated")
async def run_agent_orchestrated(request: Request):
    """Multi-Agent Role Split — Planner → N Parallel Agents (with tools) → Merger.

    요청 본문:
    {
      "prompt": "...",
      "plannerModel": "claude-opus-4-...",   // 선택
      "workerModel":  "claude-sonnet-4-6",   // 선택
      "mergerModel":  "claude-opus-4-...",   // 선택
      "maxAgents": 4,
      "awsProfile": "bedrock-gw",
      "bedrockUser": "...",
      "projectPath": "...",
      "openFile": "...", "openFileContent": "...",
      "systemPrompt": "...",                 // 선택 (스킬)
      "chatHistory": [...], "sessionId": "..."
    }

    응답: SSE
      data: {"type":"plan", "subtasks":[...]}
      data: {"type":"agent_start", "taskId":"A", ...}
      data: {"type":"agent_delta", "taskId":"A", "text":"..."}
      data: {"type":"agent_tool",  "taskId":"A", "tool":"write_file", "status":"running|done", ...}
      data: {"type":"agent_done",  "taskId":"A", "summary":"...", "toolCount":3}
      data: {"type":"merge", "report":"..."}
      data: [DONE]
    """
    body = await request.json()
    user_prompt = body.get("prompt", "")
    planner_model = body.get("plannerModel") or body.get("model") or "claude-sonnet-4-6"
    worker_model = body.get("workerModel") or body.get("model") or "claude-sonnet-4-6"
    merger_model = body.get("mergerModel") or planner_model
    max_agents = int(body.get("maxAgents", 4))
    base_sys = body.get("systemPrompt", "") or ""
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
    project_path = body.get("projectPath", "")
    open_file = body.get("openFile", "")
    open_file_content = body.get("openFileContent", "")
    # 활성 템플릿 식별자 (요구사항 5.1). 없거나 "" 이면 무템플릿(요구사항 5.2).
    template_id = body.get("templateId", "")
    is_remote = bool(body.get("isRemote", False)) or (_BRIDGE_URL and _bridge_is_remote() if _BRIDGE_URL else False)

    gw = _get_gw(aws_profile, bedrock_user)

    def _with_prefix(mid: str) -> str:
        return mid if mid.startswith("us.") or mid.startswith("eu.") else f"us.{mid}"

    planner_id = _with_prefix(planner_model)
    worker_id = _with_prefix(worker_model)
    merger_id = _with_prefix(merger_model)

    # ── Known Claude model IDs (latest gen first) — 모듈 스코프에서 import ──
    # 게이트웨이가 활성화한 모델만 호출 가능. 우선순위 순으로 시도.
    # _LATEST_OPUS_IDS / _LATEST_SONNET_IDS / _LATEST_HAIKU_IDS는 파일 상단 정의됨.
    _KNOWN_OPUS = _LATEST_OPUS_IDS
    _KNOWN_SONNET = _LATEST_SONNET_IDS
    _KNOWN_HAIKU = _LATEST_HAIKU_IDS

    def _find_model_by_keywords(keywords):
        """keywords 리스트 중 첫 번째 매칭되는 모델의 호출 가능 ID 반환."""
        # 먼저 사용자 선택 worker_model이 키워드와 일치하면 그것 사용
        for kw in keywords:
            if kw.lower() in (worker_model or "").lower():
                return worker_model
        # 알려진 모델 후보를 순회
        for kw in keywords:
            for known in _KNOWN_OPUS + _KNOWN_SONNET + _KNOWN_HAIKU:
                if kw.lower() in known.lower():
                    return known
        return None

    # 사용자가 비-Claude(예: Llama/DeepSeek) 채팅 모델을 선택해도 도구 호출 가능한 Claude로 라우팅.
    # 이렇게 해야 "PDF 만들어줘" 같은 도구 호출 작업이 실패하지 않음.
    def _is_tool_capable_chat_model(model_id: str) -> bool:
        """chat 모델이 Bedrock toolConfig를 안정적으로 지원하는지 확인."""
        if not model_id:
            return False
        mid = model_id.lower()
        # Claude는 모두 도구 호출 안정적
        if "claude" in mid:
            return True
        # Mistral Large/Pixtral도 toolConfig 지원
        if "mistral-large" in mid or "pixtral" in mid:
            return True
        # Nova Pro도 도구 호출 가능 (Nova Lite/Micro는 제한적)
        if "nova-pro" in mid:
            return True
        # 그 외 (Llama, DeepSeek-R1, Cohere Command 등) — 도구 호출 미지원/불안정
        return False

    # 사용자가 도구 호출 미지원 모델을 선택했으면 worker_id를 자동으로 Claude Sonnet으로 교체
    if not _is_tool_capable_chat_model(worker_id):
        original = worker_id
        sonnet_pref = _find_model_by_keywords(["claude-sonnet-4-6", "claude-sonnet-4", "claude-sonnet"])
        if sonnet_pref:
            worker_id = _with_prefix(sonnet_pref)
            print(f"[Orchestrator] worker 자동 변경 — {original} → {worker_id} (도구 호출 안정성)")

    # Planner도 도구 호출/JSON 출력이 안정적이어야 함 — 비-Claude면 Opus로 대체.
    # Specialized Router 사용: planner task는 항상 Opus 우선.
    planner_id = _specialized_model_for_task(
        "planner", planner_model,
        aws_profile=aws_profile, bedrock_user=bedrock_user,
    )
    merger_id = _specialized_model_for_task(
        "merger", merger_model,
        aws_profile=aws_profile, bedrock_user=bedrock_user,
    )

    # 오케스트레이터 planner/merger는 gw.converse(Bedrock)로 호출되므로 OpenAI를
    # 쓸 수 없다. 사용자가 OpenAI를 골라 보존됐더라도 Claude로 안전 대체한다.
    if is_openai_model(planner_id):
        _c = _find_model_by_keywords(["claude-opus", "claude-sonnet-4-6", "claude-sonnet"])
        if _c:
            print(f"[Orchestrator] OpenAI planner({planner_id}) → Claude({_c}) 안전 라우팅")
            planner_id = _with_prefix(_c)
    if is_openai_model(merger_id):
        _c = _find_model_by_keywords(["claude-opus", "claude-sonnet-4-6", "claude-sonnet"])
        if _c:
            print(f"[Orchestrator] OpenAI merger({merger_id}) → Claude({_c}) 안전 라우팅")
            merger_id = _with_prefix(_c)
    if is_openai_model(worker_id):
        _c = _find_model_by_keywords(["claude-sonnet-4-6", "claude-sonnet", "claude-opus"])
        if _c:
            print(f"[Orchestrator] OpenAI worker({worker_id}) → Claude({_c}) 안전 라우팅")
            worker_id = _with_prefix(_c)

    # RAG (코드 관련 질문만)
    if project_path and _is_code_related(user_prompt):
        try:
            from ai_engine.rag.context_builder import build_system_prompt
            base_sys = build_system_prompt(
                project_path=project_path, query=user_prompt,
                open_file=open_file, open_file_content=open_file_content,
                base_system_prompt=base_sys,
                aws_profile=aws_profile, bedrock_user=bedrock_user, gateway_client=gw,
            )
        except Exception as e:
            print(f"[Orchestrator] RAG 실패(무시): {e}")

    # 활성 템플릿 컨텍스트를 시스템 프롬프트에 주입 — 모델이 "템플릿 분석/반영" 요청을
    # 이해하도록 (pptx-template-styling 요구사항 5.1). 무템플릿이면 빈 문자열이라 무영향.
    _tpl_ctx = _active_template_prompt_context(template_id)
    if _tpl_ctx:
        base_sys = (base_sys or "") + _tpl_ctx

    async def orchestrated_stream():
        emit_queue: asyncio.Queue = asyncio.Queue()

        async def pipeline():
            # 자동 모델 라우팅 알림 — worker_model과 worker_id가 다르면 사용자에게 알림
            if worker_model and worker_id and worker_model.lower() not in worker_id.lower():
                await emit_queue.put({
                    "type": "model_routing",
                    "original": worker_model,
                    "routedTo": worker_id,
                    "reason": "도구 호출 안정성을 위해 Claude로 자동 라우팅됨",
                })

            # 1) Planner
            try:
                plan = await _orchestrator_plan(gw, planner_id, user_prompt, base_sys, max_agents)
            except Exception as e:
                await emit_queue.put({"type": "error", "stage": "planner", "message": str(e)})
                await emit_queue.put({"type": "__END__"})
                return

            # === Planner 결과 후처리 (Change 2) ===
            # 사용자가 명시한 형식이 있으면 plan을 검증/보강한다.
            # - 사용자가 요청하지 않은 형식의 subtask 제거 (단, 중립 작업은 유지)
            # - 사용자가 요청했지만 plan에 누락된 형식은 synthetic subtask로 강제 추가
            required_formats = _extract_required_formats(user_prompt)
            if required_formats and isinstance(plan.get("subtasks"), list):
                before_count = len(plan["subtasks"])
                plan["subtasks"] = [
                    st for st in plan["subtasks"]
                    if _subtask_format_matches(st, required_formats) or _subtask_is_neutral(st)
                ]
                have_formats = {f for f in (_subtask_primary_format(st) for st in plan["subtasks"]) if f}
                missing_formats = required_formats - have_formats
                for fmt in sorted(missing_formats):
                    if fmt in _REQUIRED_FORMAT_TOOL:
                        plan["subtasks"].append(_synthetic_subtask(fmt, user_prompt))
                if missing_formats:
                    try:
                        await emit_queue.put({
                            "type": "model_routing",
                            "original": "planner",
                            "routedTo": "format-enforcer",
                            "reason": f"누락 형식 자동 보강: {sorted(missing_formats)}",
                        })
                    except Exception:
                        pass
                # max_agents 한도 내에서 재정렬: synthetic이 truncate되지 않도록
                # synthetic을 앞쪽에 배치하면 잘리는 것을 방지.
                synthetic_ids = {f"FORCE-{f.upper()}" for f in _REQUIRED_FORMAT_TOOL}
                synthetic_first = [st for st in plan["subtasks"] if st.get("id") in synthetic_ids]
                rest = [st for st in plan["subtasks"] if st.get("id") not in synthetic_ids]
                plan["subtasks"] = synthetic_first + rest
                print(
                    f"[Planner/Post] required={sorted(required_formats)}, "
                    f"before={before_count}, after={len(plan['subtasks'])}, "
                    f"added_synthetic={sorted(missing_formats)}"
                )

            subtasks = plan["subtasks"][:max_agents]
            await emit_queue.put({"type": "plan", "subtasks": subtasks})

            # primary_tool + 작업 본문 기반 자동 모델 라우팅 — Specialized Matrix 사용.
            # 도구 호출 필요 시 Claude로 라우팅(안정성), 그 외엔 task별로 다양한 Bedrock
            # 모델 활용:
            #   - reasoning/math/logic 키워드 → DeepSeek-R1 / Qwen3
            #   - 한국어/번역/긴 문서 → Qwen3-235B (다국어 강함)
            #   - 분석/리뷰/요약 → doc_analysis (Qwen3 / Sonnet)
            #   - 단순 채팅/창작 → general_chat (Nova Lite / Gemma / Llama / Haiku)
            #   - file_generation → Sonnet (도구 호출)
            #   - code/refactor → Opus
            _file_gen_tools = {"generate_pdf", "generate_pptx", "generate_xlsx",
                               "generate_docx", "generate_image", "edit_image"}
            _code_tools = {"code", "code_analysis", "refactor"}
            # 서브태스크 본문을 보고 task type을 추정하는 키워드 매핑
            _task_keyword_map = [
                ("reasoning", ("수학", "논리", "추론", "증명", "math", "reason", "logic", "proof", "solve")),
                ("translation", ("번역", "translate", "translation", "한영", "영한", "다국어")),
                ("long_context", ("긴 문서", "장문", "전체 분석", "log analysis", "long context", "전수")),
                ("doc_analysis", ("분석", "리뷰", "검토", "평가", "비교", "analyze", "review", "compare")),
                ("korean", ("한국어", "한글", "korean")),
            ]

            def _pick_worker(st_dict):
                pt = (st_dict.get("primary_tool") or "").lower()
                # 도구 필요 작업은 안정성 위해 Claude
                if pt in _file_gen_tools:
                    task_type = "file_generation"
                elif pt in _code_tools:
                    task_type = "code"
                else:
                    # 본문/타이틀에서 task 키워드 매칭 — 비-Claude 모델 활용 기회
                    text = " ".join(str(st_dict.get(k, "") or "") for k in
                                    ("title", "description", "role")).lower()
                    matched = None
                    for tt, keywords in _task_keyword_map:
                        if any(kw in text for kw in keywords):
                            matched = tt
                            break
                    task_type = matched or "general_chat"
                picked = _specialized_model_for_task(
                    task_type, worker_model,
                    aws_profile=aws_profile, bedrock_user=bedrock_user,
                )
                # 오케스트레이터 워커 루프는 Bedrock 네이티브 tool_config 스트림이라
                # OpenAI 모델을 호출하지 못한다. 사용자가 OpenAI를 골라도 워커는
                # 안정적인 Claude로 라우팅한다(단일/병렬/합의에서는 OpenAI 그대로 사용됨).
                if is_openai_model(picked):
                    _claude_worker = _specialized_model_for_task(
                        task_type, "", aws_profile=aws_profile, bedrock_user=bedrock_user)
                    if _claude_worker and not is_openai_model(_claude_worker):
                        print(f"[Orchestrator] OpenAI 워커({picked}) → Claude({_claude_worker}) 안전 라우팅")
                        picked = _claude_worker
                return picked

            # 2) Hierarchical or Parallel Agents
            # 팀(team) 필드가 2개 이상 다른 값으로 있으면 hierarchical 모드 활성:
            # - 같은 team끼리 묶어서 순차 실행 (team A → team B)
            # - 같은 team 안에서는 병렬 실행
            # - 이전 team의 결과(verifiedFiles + summary)가 다음 team 컨텍스트로 전달
            # 이 모드는 기존 flat 흐름과 결과 동일 — 단지 정보 흐름 개선만.
            teams_set = set()
            for st in subtasks:
                t = (st.get("team") or "").strip().lower()
                if t:
                    teams_set.add(t)

            use_hierarchical = len(teams_set) >= 2

            # 전체 오케스트레이션 timeout — 환경변수로 조정 (기본 60분)
            total_timeout = float(os.environ.get("AE_ORCH_TOTAL_TIMEOUT", "3600"))

            agent_results = []
            if use_hierarchical:
                print(f"[Orchestrator] Hierarchical 모드 — {len(teams_set)}개 팀 ({sorted(teams_set)})")
                await emit_queue.put({
                    "type": "hierarchical_info",
                    "teams": sorted(teams_set),
                    "totalAgents": len(subtasks),
                })
                # 팀 우선순위: research/analysis 먼저, 그 다음 coding/media/writing
                team_priority = {"research": 1, "analysis": 1, "coding": 2,
                                 "media": 3, "writing": 4}
                ordered_teams = sorted(teams_set, key=lambda t: team_priority.get(t, 5))
                team_summaries = {}  # team → 직전 결과 요약 (다음 team 컨텍스트)
                remaining_timeout = total_timeout

                for team in ordered_teams:
                    team_subtasks = [st for st in subtasks
                                     if (st.get("team") or "").strip().lower() == team]
                    # 이전 팀 결과를 컨텍스트로 주입
                    team_base_sys = base_sys
                    if team_summaries:
                        prev_ctx = "\n\n".join(
                            f"## {prev_team} 팀 결과 (참고용)\n{summary[:1500]}"
                            for prev_team, summary in team_summaries.items()
                        )
                        team_base_sys = (base_sys or "") + "\n\n[이전 팀 결과]\n" + prev_ctx

                    team_tasks = [
                        asyncio.create_task(
                            _orchestrator_run_agent(gw, _pick_worker(st), st, project_path,
                                                    team_base_sys, emit_queue,
                                                    aws_profile=aws_profile, bedrock_user=bedrock_user,
                                                    is_remote=is_remote, template_id=template_id)
                        )
                        for st in team_subtasks
                    ]
                    # 각 팀 timeout — 남은 전체 timeout / 남은 팀 수
                    team_timeout = max(60, remaining_timeout / max(1, len(ordered_teams) - len(team_summaries)))
                    team_start = asyncio.get_event_loop().time()
                    try:
                        team_raw = await asyncio.wait_for(
                            asyncio.gather(*team_tasks, return_exceptions=True),
                            timeout=team_timeout,
                        )
                    except asyncio.TimeoutError:
                        team_raw = []
                        for t in team_tasks:
                            if t.done():
                                try: team_raw.append(t.result())
                                except Exception as ex: team_raw.append(ex)
                            else:
                                t.cancel()
                                team_raw.append(asyncio.TimeoutError(f"team {team} timeout"))
                    except Exception as e:
                        team_raw = [Exception(str(e))] * len(team_subtasks)
                    elapsed = asyncio.get_event_loop().time() - team_start
                    remaining_timeout = max(60, remaining_timeout - elapsed)

                    # 팀 결과 정규화 + summary 수집
                    summary_parts = []
                    for st, r in zip(team_subtasks, team_raw):
                        tid = st.get("id", "?")
                        if isinstance(r, Exception):
                            err_msg = f"{type(r).__name__}: {r}"
                            await emit_queue.put({"type": "agent_error", "taskId": tid, "error": err_msg})
                            agent_results.append({
                                "taskId": tid, "role": st.get("role", "Worker"),
                                "title": st.get("title", ""), "status": "error",
                                "summary": err_msg, "tools": [], "verifiedFiles": [],
                                "team": team,
                            })
                        elif isinstance(r, dict):
                            r["team"] = team
                            agent_results.append(r)
                            if r.get("status") == "done":
                                summary_parts.append(
                                    f"- [{tid}] {st.get('title', '')}: "
                                    f"{(r.get('summary') or '')[:300]}"
                                )
                        else:
                            agent_results.append({
                                "taskId": tid, "role": st.get("role", "Worker"),
                                "title": st.get("title", ""), "status": "error",
                                "summary": "에이전트가 결과를 반환하지 않음",
                                "tools": [], "verifiedFiles": [], "team": team,
                            })
                    if summary_parts:
                        team_summaries[team] = "\n".join(summary_parts)
            else:
                # === Flat parallel — 기존 동작 ===
                tasks = [
                    asyncio.create_task(
                        _orchestrator_run_agent(gw, _pick_worker(st), st, project_path,
                                                base_sys, emit_queue,
                                                aws_profile=aws_profile, bedrock_user=bedrock_user,
                                                is_remote=is_remote, template_id=template_id)
                    )
                    for st in subtasks
                ]
                try:
                    raw_results = await asyncio.wait_for(
                        asyncio.gather(*tasks, return_exceptions=True),
                        timeout=total_timeout,
                    )
                except asyncio.TimeoutError:
                    print(f"[Orchestrator] 전체 timeout ({total_timeout}s) — 미완료 태스크 취소")
                    raw_results = []
                    for t in tasks:
                        if t.done():
                            try: raw_results.append(t.result())
                            except Exception as ex: raw_results.append(ex)
                        else:
                            t.cancel()
                            raw_results.append(asyncio.TimeoutError(f"timeout {int(total_timeout)}s"))
                except Exception as e:
                    print(f"[Orchestrator] gather 예외: {e}")
                    raw_results = [Exception(str(e))] * len(subtasks)

                # 예외를 결과 dict로 정규화
                for st, r in zip(subtasks, raw_results):
                    tid = st.get("id", "?")
                    if isinstance(r, Exception):
                        err_msg = f"{type(r).__name__}: {r}"
                        await emit_queue.put({"type": "agent_error", "taskId": tid, "error": err_msg})
                        agent_results.append({
                            "taskId": tid, "role": st.get("role", "Worker"),
                            "title": st.get("title", ""), "status": "error",
                            "summary": err_msg, "tools": [], "verifiedFiles": [],
                        })
                    elif isinstance(r, dict):
                        agent_results.append(r)
                    else:
                        agent_results.append({
                            "taskId": tid, "role": st.get("role", "Worker"),
                            "title": st.get("title", ""), "status": "error",
                            "summary": "에이전트가 결과를 반환하지 않음",
                            "tools": [], "verifiedFiles": [],
                        })

            # === Final pipeline format coverage verification (Change 5) ===
            # 사용자가 명시한 형식 중 디스크에 0건인 형식이 있으면, merger 전에
            # synthetic agent_result로 강제 보강한다. 이 단계가 "1개 성공" 같은
            # false-positive를 차단하는 핵심.
            import time as _time_local
            if required_formats:
                all_verified_paths_lc = []
                for r in agent_results:
                    for vf in (r.get("verifiedFiles") or []):
                        if isinstance(vf, dict) and vf.get("path"):
                            all_verified_paths_lc.append(str(vf["path"]).lower())

                have_formats = set()
                for p in all_verified_paths_lc:
                    for ext in ("pdf", "pptx", "xlsx", "docx", "png", "jpg", "jpeg"):
                        if p.endswith(f".{ext}"):
                            have_formats.add("png" if ext in ("jpg", "jpeg") else ext)
                            break

                missing_cov = required_formats - have_formats
                if missing_cov:
                    await emit_queue.put({
                        "type": "model_routing",
                        "original": "verification",
                        "routedTo": "format-coverage-enforcer",
                        "reason": f"사용자 요청 형식 누락: {sorted(missing_cov)} — 추가 생성 시도",
                    })
                    worker_id_for_enrich = _specialized_model_for_task(
                        "enrich", worker_model,
                        aws_profile=aws_profile, bedrock_user=bedrock_user,
                    )
                    for fmt in sorted(missing_cov):
                        primary_tool_cov = _REQUIRED_FORMAT_TOOL.get(fmt)
                        if not primary_tool_cov:
                            continue
                        # 다른 에이전트의 성공 summary들을 cross-pollinate해서 시드 본문 생성
                        seed_text = "\n\n".join(
                            (r.get("summary") or "")[:1500]
                            for r in agent_results
                            if r.get("status") == "done"
                        ) or user_prompt
                        try:
                            enriched_cov = await _enrich_content_via_gateway(
                                gw=gw, model_id=worker_id_for_enrich,
                                primary_tool=primary_tool_cov,
                                title=user_prompt[:60] or "사용자 요청",
                                description=user_prompt,
                                final_text=seed_text,
                                project_path=project_path,
                            )
                            forced_cov = await _force_generate_from_text(
                                primary_tool=primary_tool_cov,
                                target_files=[f".generated/coverage-{fmt}-{int(_time_local.time())}.{fmt}"],
                                title=user_prompt[:60] or "사용자 요청",
                                description=user_prompt,
                                final_text=enriched_cov or user_prompt,
                                project_path=project_path,
                                aws_profile=aws_profile,
                                bedrock_user=bedrock_user,
                                template_id=template_id,  # 요구사항 5.1 — 활성 템플릿 전달
                            )
                            if forced_cov:
                                synthetic_id = f"COV-{fmt.upper()}"
                                agent_results.append({
                                    "taskId": synthetic_id,
                                    "role": f"Coverage Enforcer ({fmt.upper()})",
                                    "title": f"누락 형식 보강 — {fmt.upper()}",
                                    "status": "done",
                                    "summary": f"[형식 보강] {fmt.upper()} 파일을 추가 생성했습니다.",
                                    "tools": [{"name": f"coverage:{primary_tool_cov}",
                                               "input": {"format": fmt},
                                               "output": forced_cov[0][0]}],
                                    "verifiedFiles": [vf for _, vf in forced_cov],
                                })
                                await emit_queue.put({
                                    "type": "agent_done", "taskId": synthetic_id,
                                    "summary": f"누락 {fmt.upper()} 형식 자동 보강 완료",
                                    "toolCount": 1,
                                    "verifiedFiles": [vf["path"] for _, vf in forced_cov],
                                })
                        except Exception as _ce:
                            print(f"[Orchestrator] coverage 보강 실패 ({fmt}): {_ce}")

            # 3) Merger + 합의 교차검증(병렬·비차단·additive)
            #    병합과 교차검증을 asyncio.gather로 동시 실행 → 교차검증이 느리거나
            #    degraded여도 병합(report)은 지연/차단되지 않는다(가용성 우선).
            #    자동 ON(공수 0). 끄려면 AE_CONSENSUS_CROSSVERIFY=0.
            _cv_flag = str(os.environ.get("AE_CONSENSUS_CROSSVERIFY") or "").strip().lower()
            _cv_on = _cv_flag not in ("0", "false", "no", "off")

            async def _do_merge():
                return await _orchestrator_merge(gw, merger_id, user_prompt, agent_results, base_sys,
                                                 required_formats=required_formats)

            async def _do_crossverify():
                # 후보가 2개 미만이면 합의 자체가 무의미 → skip.
                if not _cv_on or len([r for r in agent_results if r.get("status") == "done"]) < 2:
                    return None
                try:
                    from ai_engine.rag.cross_verify import cross_verify_consensus
                    try:
                        _cv_to = float(os.environ.get("AE_CONSENSUS_CROSSVERIFY_TIMEOUT_MS", "12000")) / 1000.0
                    except (TypeError, ValueError):
                        _cv_to = 12.0
                    _rep = await cross_verify_consensus(gw, merger_id, user_prompt,
                                                        agent_results, timeout=_cv_to)
                    return _rep.as_dict()
                except Exception as _cve:
                    print(f"[Orchestrator] cross-verify 예외(비차단): {_cve}")
                    return None

            report, _cv = await asyncio.gather(_do_merge(), _do_crossverify())
            _merge_evt = {"type": "merge", "report": report, "results": agent_results}
            if _cv is not None:
                _merge_evt["crossVerify"] = _cv
                _cc = _cv.get("conflictCount", 0)
                print(f"[Orchestrator] cross-verify done — conflicts={_cc} "
                      f"degraded={_cv.get('degraded')}")
            await emit_queue.put(_merge_evt)
            await emit_queue.put({"type": "__END__"})

        pipe_task = asyncio.create_task(pipeline())

        # Heartbeat — 30초마다 클라이언트 idle timeout 방지
        last_send = asyncio.get_event_loop().time()
        try:
            while True:
                try:
                    evt = await asyncio.wait_for(emit_queue.get(), timeout=20)
                    if evt.get("type") == "__END__":
                        break
                    try:
                        _line = f"data: {json.dumps(evt, ensure_ascii=False)}\n\n"
                    except (TypeError, ValueError) as _je:
                        # 직렬화 불가 이벤트 1건 때문에 전체 스트림이 끊기지 않도록 방어.
                        _line = f"data: {json.dumps({'type': 'agent_tool', 'status': 'done', 'output': 'event-serialize-skip'}, ensure_ascii=False)}\n\n"
                        print(f"[Orchestrator] 이벤트 직렬화 실패(스킵): {str(_je)[:160]}")
                    yield _line
                    last_send = asyncio.get_event_loop().time()
                except asyncio.TimeoutError:
                    # 20초 동안 이벤트 없음 → heartbeat
                    yield f"data: {json.dumps({'heartbeat': True, 'ts': int(asyncio.get_event_loop().time())})}\n\n"
                    last_send = asyncio.get_event_loop().time()
                    # pipeline이 끝났는지 확인
                    if pipe_task.done():
                        break
        finally:
            if not pipe_task.done():
                pipe_task.cancel()
                try: await pipe_task
                except Exception: pass
            yield "data: [DONE]\n\n"

    return StreamingResponse(orchestrated_stream(), media_type="text/event-stream")


@app.post("/api/agents/run")
async def run_agent(request: Request):
    body = await request.json()
    prompt = body.get("prompt", "")
    model = body.get("model", "anthropic.claude-sonnet-4-5-20250929-v1:0")
    aws_profile = body.get("awsProfile", os.environ.get("AWS_PROFILE", "bedrock-gw"))
    bedrock_user = body.get("bedrockUser", os.environ.get("BEDROCK_USER", ""))
    try:
        gw = _get_gw(aws_profile, bedrock_user)
        result = await gw.converse(
            model_id=model,
            messages=[{"role": "user", "content": [{"text": prompt}]}],
        )
        return JSONResponse(content=result)
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)


@app.post("/api/conversation/handoff")
async def conversation_handoff(request: Request):
    """대화 인계 — 누적 컨텍스트가 임계치를 넘으면 클라이언트가 호출.

    Haiku로 이전 대화를 요약 → .generated/handoff/conversation-handoff-<ts>.md 저장.
    클라이언트는 새 탭을 자동 생성하고 요약 MD를 첫 user 메시지로 주입한다.

    이전 버전: 클라이언트는 호출하는데 서버에 endpoint가 없어 매 누적 메시지마다
    404가 콘솔에 박혔음. 이 endpoint가 그 404를 해소함.
    """
    try:
        body = await request.json()
        msgs = body.get("messages", []) or []
        session_id = body.get("sessionId", "") or ""
        project_path = body.get("projectPath", "") or ""
        aws_profile = body.get("awsProfile", "bedrock-gw") or "bedrock-gw"
        bedrock_user = body.get("bedrockUser", "") or ""

        # 보존 디렉토리
        local_root = _resolve_local_root(project_path)
        handoff_dir = os.path.join(local_root, ".generated", "handoff")
        os.makedirs(handoff_dir, exist_ok=True)
        ts = int(__import__("time").time() * 1000)
        out_name = f"conversation-handoff-{ts}.md"
        out_path = os.path.join(handoff_dir, out_name)
        rel_path = f".generated/handoff/{out_name}"

        # 본문 — Haiku로 요약 시도, 실패 시 deterministic 요약
        flat_text = "\n\n".join(
            f"## {m.get('role','user').upper()}\n{(m.get('content') or '')[:3000]}"
            for m in msgs[-30:]
        )
        summary_prompt = (
            "당신은 대화 인계 요약가입니다. 다음 대화의 핵심 작업, 결정 사항, "
            "미완료 항목을 5~10개 글머리표로 정리하세요. 새 대화창에서 이어서 작업할 "
            "수 있도록 충분한 컨텍스트를 포함합니다.\n\n"
            f"--- 대화 ---\n{flat_text[:60000]}"
        )
        summary_text = ""
        try:
            gw = _get_gw(aws_profile, bedrock_user)
            haiku = _specialized_model_for_task(
                "summarize", "", aws_profile=aws_profile, bedrock_user=bedrock_user
            ) or "anthropic.claude-haiku-4-5-20251001-v1:0"
            haiku = _resolve_callable_model_id(haiku, aws_profile, bedrock_user)
            result = await asyncio.wait_for(
                gw.converse(
                    haiku,
                    [{"role": "user", "content": [{"text": summary_prompt}]}],
                    "당신은 간결한 한국어 요약 전문가입니다.",
                ),
                timeout=60.0,
            )
            if result and result.get("decision") == "ALLOW":
                content = result.get("output", {}).get("message", {}).get("content", [])
                summary_text = "\n".join(c.get("text", "") for c in content if "text" in c)
        except Exception as e:
            print(f"[Handoff] LLM 요약 실패 — deterministic fallback: {e}")

        if not summary_text or len(summary_text) < 50:
            # Deterministic fallback — 메시지 헤드 + 마지막 user 발화
            heads = []
            for m in msgs[:5]:
                snippet = (m.get("content") or "")[:300].replace("\n", " ")
                heads.append(f"- [{m.get('role','user')}] {snippet}")
            tail = msgs[-1] if msgs else {}
            tail_snippet = (tail.get("content") or "")[:1000]
            summary_text = (
                "## 대화 인계 요약 (자동 생성)\n\n"
                "### 시작 부분\n" + "\n".join(heads) + "\n\n"
                f"### 마지막 사용자 메시지\n{tail_snippet}\n\n"
                "### 인계 노트\n- LLM 요약을 사용할 수 없어 deterministic fallback으로 생성됨\n"
                f"- 총 {len(msgs)}개 메시지\n"
            )

        md_content = (
            f"# 대화 인계 — {datetime.utcnow().isoformat()}Z\n\n"
            f"세션 ID: `{session_id or '-'}`\n\n"
            f"메시지 수: {len(msgs)}\n\n"
            f"---\n\n{summary_text}\n"
        )
        try:
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(md_content)
        except Exception as e:
            return JSONResponse(content={"error": f"파일 저장 실패: {e}"}, status_code=500)

        return JSONResponse(content={
            "ok": True,
            "path": rel_path,
            "absPath": out_path,
            "content": summary_text,
            "messageCount": len(msgs),
        })
    except Exception as e:
        return JSONResponse(content={"error": str(e)}, status_code=500)


_quota_cache = {"used_krw": 0, "remaining_krw": 0, "limit_krw": 0, "last_updated": "", "user": "", "fetching": False}

@app.get("/api/quota")
async def get_quota(request: Request):
    profile = request.query_params.get("profile", os.environ.get("AWS_PROFILE", "default"))
    user = request.query_params.get("user", "")
    # 첫 호출이면 백그라운드에서 quota 조회 시작 (즉시 응답 반환)
    if _quota_cache["remaining_krw"] == 0 and user and not _quota_cache["fetching"]:
        _quota_cache["fetching"] = True
        asyncio.create_task(_fetch_quota_background(profile, user))
    return {
        "user": _quota_cache["user"] or user,
        "used_krw": round(_quota_cache["used_krw"], 2),
        "remaining_krw": round(_quota_cache["remaining_krw"], 2),
        "limit_krw": round(_quota_cache["limit_krw"], 2),
        "last_updated": _quota_cache["last_updated"],
    }


async def _fetch_quota_background(profile, user):
    """백그라운드에서 Gateway 호출하여 quota 정보 캐시."""
    try:
        gw = _get_gw(profile, user)
        print(f"[Quota] 백그라운드 quota 조회 시작...")
        # Gateway /converse 직접 호출 — maxTokens:1로 최소 비용
        # us. prefix haiku 4.5 우선, 실패 시 haiku 3 fallback.
        # 첫 항목은 라우터로 — 항상 최신 Haiku 사용.
        primary_haiku = _specialized_model_for_task(
            "simple_qa", None,
            aws_profile=profile, bedrock_user=user,
        )
        quota_models = [
            primary_haiku,
            "us.anthropic.claude-haiku-4-5-20251001-v1:0",
            "us.anthropic.claude-3-haiku-20240307-v1:0",
        ]
        # 중복 제거
        seen = set()
        quota_models = [m for m in quota_models if m and not (m in seen or seen.add(m))]
        for mid in quota_models:
            try:
                result = await asyncio.wait_for(
                    gw.converse_quota_only(
                        model_id=mid,
                        messages=[{"role": "user", "content": [{"text": "hi"}]}],
                    ),
                    timeout=15
                )
                print(f"[Quota] {mid}: decision={result.get('decision')}, remaining_quota={result.get('remaining_quota')}, error={result.get('error', '')[:100]}")
                _extract_quota(result, user)
                if _quota_cache["remaining_krw"] > 0:
                    print(f"[Quota] 성공! remaining={_quota_cache['remaining_krw']}")
                    return
            except Exception as e:
                print(f"[Quota] {mid} 실패: {e}")
                continue
        print(f"[Quota] 모든 모델 실패 — 첫 채팅 후 자동 갱신")
    except Exception as e:
        print(f"[Quota] 백그라운드 조회 실패: {e}")
    finally:
        _quota_cache["fetching"] = False


def _extract_quota(result, user=""):
    """Gateway 응답에서 quota 정보를 추출하여 캐시에 저장."""
    rq = result.get("remaining_quota", {})
    if rq:
        # 다양한 키 이름 대응
        cost_val = 0
        if isinstance(rq, (int, float)):
            cost_val = rq
        elif isinstance(rq, dict):
            cost_val = rq.get("cost_krw") or rq.get("remaining_cost_krw") or rq.get("remaining_krw") or rq.get("remaining") or 0
            if not cost_val:
                for v in rq.values():
                    if isinstance(v, (int, float)) and v > 0:
                        cost_val = v
                        break
        if cost_val > 0:
            _quota_cache["remaining_krw"] = cost_val
            _quota_cache["limit_krw"] = cost_val + _quota_cache["used_krw"]
            _quota_cache["user"] = user
            _quota_cache["last_updated"] = datetime.utcnow().isoformat()
            print(f"[Quota] 캐시 갱신 성공: remaining={cost_val}")
    if result.get("estimated_cost_krw"):
        _quota_cache["used_krw"] += result["estimated_cost_krw"]
        if _quota_cache["remaining_krw"] > 0:
            _quota_cache["limit_krw"] = _quota_cache["remaining_krw"] + _quota_cache["used_krw"]


# ===========================================================================
# 템플릿 엔드포인트 5종 — PPTX 템플릿 스타일링
# 설계 §구성요소 9 (FastAPI 엔드포인트 계약), 요구사항 1.5/2.7/3.5/5.3/8.1/8.8/9.3
# ===========================================================================
#
# Template_Manager(ai_engine/template_manager.py)의 register/list/get/
# get_style_profile/delete 결과를 JSONResponse로 변환해 노출한다. Template_Manager는
# 어떤 입력/오류에도 예외를 밖으로 던지지 않고 항상 dict를 반환하므로(폴백 격리 원칙),
# 여기서는 그 결과를 그대로 직렬화하고 에러 이름에 맞는 HTTP status code만 매핑한다
# (설계 §에러 이름 ↔ 요구사항 매핑 표). template_manager는 핸들러 내부에서 dual-path로
# 지연 import 하므로(slide_templates 관례와 동일) 모듈 import 자체는 항상 성공한다.

# 에러 문자열 → HTTP status code 매핑(설계 §에러 이름 ↔ 요구사항 매핑 표).
_TEMPLATE_ERROR_STATUS = {
    "invalid-name": 400,            # 1.2, 1.7
    "invalid-template": 400,        # 1.3
    "template-too-large": 400,      # 1.4
    "invalid-template-id": 400,     # 2.7
    "invalid-json": 400,            # 4.5
    "invalid-style-profile": 400,   # 4.6
    "invalid-color": 400,           # 4.7
    "duplicate-name": 409,          # 1.6 (충돌)
    "template-not-found": 404,      # 5.4, 8.x
    "no-storage-root": 500,         # 2.4
    "template-store-write-failed": 500,  # 2.6
    "template-delete-failed": 500,  # 8.12
    "missing-dep": 503,             # 9.3 (의존성 부재)
}


def _load_template_manager():
    """template_manager 모듈을 dual-path로 지연 import 한다(slide_templates 관례와 동일).

    repo 루트 실행 시 `ai_engine.template_manager`, ai_engine/ 내부 실행 시
    `template_manager`. 둘 다 실패하면 None을 반환하고, 라우트가 missing-dep(503)으로
    처리한다. template_manager는 최상단에서 json/os만 import 하므로 정상 환경에서는
    항상 성공한다(python-pptx 등 무거운 의존성은 등록/추출 시점에 지연 import 됨).
    """
    try:
        from ai_engine import template_manager
        return template_manager
    except ImportError:
        try:
            import template_manager  # type: ignore  # alt path when running from ai_engine/
            return template_manager
        except ImportError:
            return None


def _template_response(result):
    """Template_Manager 결과 dict를 적절한 status code의 JSONResponse로 변환한다.

    성공(에러 없음)은 200, `error` 필드가 있으면 _TEMPLATE_ERROR_STATUS 매핑을 따른다
    (미등록 에러는 400). 응답 body는 Template_Manager가 반환한 dict를 그대로 담아
    에러 이름·부가 필드(lib/hint/maxBytes/allowed 등)를 보존한다.
    """
    status = 200
    if isinstance(result, dict):
        err = result.get("error")
        if err:
            status = _TEMPLATE_ERROR_STATUS.get(err, 400)
    return JSONResponse(content=result, status_code=status)


def _template_manager_unavailable():
    """template_manager 모듈을 import 할 수 없을 때의 missing-dep 응답(503)."""
    return JSONResponse(
        content={
            "error": "missing-dep",
            "lib": "template_manager",
            "hint": "ai_engine/template_manager.py를 찾을 수 없습니다",
        },
        status_code=503,
    )


def _active_template_prompt_context(template_id: str) -> str:
    """활성 템플릿의 스타일·구조 정보를 LLM 시스템 프롬프트에 주입할 텍스트로 만든다.

    채팅/오케스트레이터 모델이 "등록된 템플릿을 분석해서 재구성해" 같은 요청을 이해하려면,
    활성 템플릿이 무엇이고 어떤 색/폰트/레이아웃을 갖는지 *프롬프트로* 알아야 한다. 생성
    단계(_tool_generate_pptx)는 base.pptx를 상속하지만, 그 사실을 모델은 알 수 없기 때문이다.

    template_id가 없거나("" → 무템플릿) 해석에 실패하면 빈 문자열을 반환한다(기존 동작 보존).
    어떤 예외에도 빈 문자열로 폴백한다(폴백 격리 — gateway.md 위반 없음: 신규 LLM 호출 없이
    로컬 파일 파싱 결과를 데이터로만 주입).

    Returns:
        주입할 컨텍스트 문자열(앞에 개행 포함), 또는 무템플릿/실패 시 "".
    """
    if not template_id:
        return ""
    try:
        tm = _load_template_manager()
        if tm is None:
            return ""
        store_root = tm.resolve_template_store_root()
        if not store_root:
            return ""
        meta = tm.get_template(template_id, store_root)
        if not isinstance(meta, dict) or meta.get("error"):
            return ""
        name = meta.get("name") or "(이름 없음)"
        layout_count = meta.get("layoutCount")
        sp = meta.get("styleProfile") or {}
        layouts = meta.get("layouts") or []
        role_counts = meta.get("roleCounts") or {}
        # styleProfile이 dict가 아니면(손상 등) 색/폰트 줄은 생략하고 이름만이라도 알린다.
        lines = [
            "",
            "[활성 PPTX 템플릿]",
            f"사용자가 다음 PowerPoint 템플릿을 활성 템플릿으로 선택했습니다: \"{name}\".",
            "이후 생성되는 PPTX는 이 템플릿의 슬라이드 마스터·레이아웃·테마를 자동 상속하며,",
            "아래 색/폰트 토큰이 배경·다이어그램·차트에 적용됩니다. 슬라이드를 구성할 때 이",
            "스타일과 조화를 이루도록 색·폰트·톤을 맞추고, 사용자가 '템플릿을 분석/반영해 달라'고",
            "요청하면 아래 정보를 근거로 답하세요.",
        ]
        if isinstance(sp, dict) and sp:
            def _g(k):
                v = sp.get(k)
                return v if isinstance(v, str) and v else "—"
            lines.append(
                "- 색상 팔레트: "
                f"주 {_g('primaryColor')}, 보조 {_g('secondaryColor')}, "
                f"강조 {_g('accentColor')}, 텍스트 {_g('textColor')}, 배경 {_g('backgroundColor')}"
            )
            lines.append(
                f"- 글꼴: 제목 \"{_g('headingFont')}\", 본문 \"{_g('bodyFont')}\""
            )
        if isinstance(layout_count, int) and layout_count > 0:
            lines.append(f"- 사용 가능한 슬라이드 레이아웃 수: {layout_count}개")

        # 레이아웃 구조 — 모델이 콘텐츠를 적합한 레이아웃에 매핑하도록 역할별로 제시.
        if isinstance(layouts, list) and layouts:
            _role_label = {
                "title": "표지",
                "two-column": "2단/비교",
                "section": "섹션 구분",
                "content": "제목+내용",
                "blank": "빈 화면",
                "other": "기타",
            }
            # 역할 요약(있으면)
            if isinstance(role_counts, dict) and role_counts:
                summary = ", ".join(
                    f"{_role_label.get(r, r)} {c}개"
                    for r, c in role_counts.items() if c
                )
                if summary:
                    lines.append(f"- 레이아웃 구성 요약: {summary}")
            lines.append("- 슬라이드 레이아웃 목록(역할 → 이름):")
            for lay in layouts[:24]:  # 과도한 토큰 방지 상한
                if not isinstance(lay, dict):
                    continue
                r = _role_label.get(lay.get("role"), lay.get("role") or "기타")
                nm = lay.get("name") or "(이름 없음)"
                lines.append(f"  · {r} → \"{nm}\"")
            lines.append(
                "콘텐츠를 슬라이드로 나눌 때 각 슬라이드의 성격에 맞는 레이아웃 역할을 고르세요: "
                "표지/마무리는 '표지', 두 항목 대조·장단점은 '2단/비교', 챕터 전환은 '섹션 구분', "
                "일반 본문은 '제목+내용'. 각 슬라이드 객체의 layout 필드를 "
                "'title'/'two-column'/'content' 중 적절한 값으로 지정하면 그 역할의 레이아웃에 매핑됩니다."
            )
        lines.append("")
        return "\n".join(lines)
    except Exception as exc:  # 폴백 격리 — 어떤 실패도 무템플릿처럼 동작
        try:
            print(f"[ActiveTemplate] 프롬프트 컨텍스트 생성 실패(무시): {str(exc)[:200]}")
        except Exception:
            pass
        return ""


@app.post("/api/media/pptx-render")
async def api_media_pptx_render(request: Request):
    """생성된 .pptx를 앱 미리보기용으로 파싱한다.

    응답: {"slides": [{"title": str, "bullets": [str], "images": [dataURL]}]}.
    임베드된 Picture(풀블리드 Vertex/HTML 이미지 포함)는 base64 data URL로 인라인해
    미리보기에서 실제 시각 품질을 확인할 수 있게 한다. python-pptx만 사용(외부 렌더러 불필요).
    """
    try:
        body = await request.json()
    except Exception:
        body = {}
    rel = (body.get("path") or "").strip()
    b64 = (body.get("base64") or "").strip()
    if not rel and not b64:
        return JSONResponse({"error": "path 또는 base64가 필요합니다."}, status_code=400)
    # 경로 해석 — 절대경로 우선, 아니면 userData/.generated 및 repo .generated 후보 탐색.
    cand = rel
    if rel and not os.path.isabs(cand):
        roots = []
        try:
            roots.append(_resolve_local_root(""))
        except Exception:
            pass
        roots.append(os.getcwd())
        for r in roots:
            c = os.path.join(r, rel)
            if os.path.isfile(c):
                cand = c
                break
    # path가 안 잡히면 base64 폴백 — 프론트가 이미 파일 내용을 갖고 있으므로
    # 경로 해석 실패(원격/슬러그 불일치 등)에도 미리보기를 보장한다.
    _tmp_preview = ""
    if (not cand or not os.path.isfile(cand)) and b64:
        try:
            import base64 as _b64dec, tempfile as _tfp
            _raw = _b64dec.b64decode(b64)
            _tf = _tfp.NamedTemporaryFile(prefix="ae_pptx_preview_", suffix=".pptx", delete=False)
            _tf.write(_raw)
            _tf.close()
            cand = _tf.name
            _tmp_preview = _tf.name
        except Exception as _be:
            return JSONResponse({"error": f"base64 디코드 실패: {str(_be)[:120]}"}, status_code=400)
    if not cand or not os.path.isfile(cand):
        return JSONResponse({"error": f"파일을 찾을 수 없습니다: {rel}"}, status_code=404)
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        import base64 as _b64
    except Exception as e:
        return JSONResponse({"error": f"pptx 파서 로드 실패: {str(e)[:160]}"}, status_code=500)
    try:
        prs = Presentation(cand)
    except Exception as e:
        return JSONResponse({"error": f"pptx 열기 실패: {str(e)[:160]}"}, status_code=500)

    def _mime(blob: bytes, ext: str) -> str:
        if blob[:8] == b"\x89PNG\r\n\x1a\n":
            return "image/png"
        if blob[:3] == b"\xff\xd8\xff":
            return "image/jpeg"
        e = (ext or "").lower().lstrip(".")
        return f"image/{e}" if e else "image/png"

    slides_out = []
    MAX_IMG_BYTES = 4_000_000  # 미리보기 응답 비대화 방지(>4MB 이미지는 생략)
    for sl in prs.slides:
        title = ""
        bullets = []
        images = []
        for sh in sl.shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    blob = sh.image.blob
                    if blob and len(blob) <= MAX_IMG_BYTES:
                        mime = _mime(blob, getattr(sh.image, "ext", "png"))
                        images.append(f"data:{mime};base64," + _b64.b64encode(blob).decode("ascii"))
                    continue
            except Exception:
                pass
            if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if not txt:
                    continue
                is_title = False
                try:
                    if sh == sl.shapes.title:
                        is_title = True
                except Exception:
                    pass
                if is_title and not title:
                    title = txt.splitlines()[0]
                else:
                    for line in txt.splitlines():
                        line = line.strip()
                        if line:
                            bullets.append(line)
        slides_out.append({"title": title, "bullets": bullets, "images": images})
    return JSONResponse({"slides": slides_out})


@app.post("/api/templates")
async def api_register_template(request: Request):
    """템플릿 등록 (요구사항 1.5, 9.3). body {filePath, name} → register_template()."""
    tm = _load_template_manager()
    if tm is None:
        return _template_manager_unavailable()
    try:
        body = await request.json()
    except Exception:
        body = {}
    if not isinstance(body, dict):
        body = {}
    file_path = body.get("filePath", "")
    name = body.get("name", "")
    result = tm.register_template(file_path, name)
    return _template_response(result)


@app.get("/api/templates")
async def api_list_templates(request: Request):
    """템플릿 목록 (요구사항 8.1). {"templates": [...]} createdAt 내림차순, 최대 200개."""
    tm = _load_template_manager()
    if tm is None:
        # 모듈 부재 시에도 목록은 빈 배열로 응답(목록 조회는 오류가 아님).
        return JSONResponse(content={"templates": []})
    return _template_response(tm.list_templates())


@app.get("/api/templates/{template_id}")
async def api_get_template(template_id: str):
    """단건 조회 (요구사항 5.3, 2.7). get_template() 결과를 그대로 반환."""
    tm = _load_template_manager()
    if tm is None:
        return _template_manager_unavailable()
    return _template_response(tm.get_template(template_id))


@app.get("/api/templates/{template_id}/style-profile")
async def api_get_template_style_profile(template_id: str):
    """Style_Profile 조회 (요구사항 3.5). 저장된 style_profile.json을 매 호출 동일하게 반환.

    store_root를 먼저 해석한 뒤 get_style_profile(template_id, store_root)을 호출한다.
    store_root 결정 불가 시 no-storage-root, 그 외에는 manager가 반환하는 parsed profile
    dict 또는 {"error": invalid-template-id | template-not-found | invalid-json}.
    """
    tm = _load_template_manager()
    if tm is None:
        return _template_manager_unavailable()
    store_root = tm.resolve_template_store_root()
    if not store_root:
        return _template_response({"error": "no-storage-root"})
    result = tm.get_style_profile(template_id, store_root)
    return _template_response(result)


@app.delete("/api/templates/{template_id}")
async def api_delete_template(template_id: str):
    """템플릿 삭제 (요구사항 8.8, 2.7). delete_template() 결과를 그대로 반환."""
    tm = _load_template_manager()
    if tm is None:
        return _template_manager_unavailable()
    return _template_response(tm.delete_template(template_id))
