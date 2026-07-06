"""Template_Manager — PPTX 프레젠테이션 템플릿의 등록·조회·삭제·저장 격리.

이 모듈은 사용자가 등록한 PowerPoint 템플릿(`.pptx`)을 `userData/templates/{templateId}/`
하위에만 안전하게 보관하고, 그 테마에서 Style_Profile(색/폰트 토큰)을 파생하는 백엔드
구성요소다. python-pptx 의존성은 등록/추출 시점에 지연 import 하여 미설치 시
"missing-dep"를 반환한다(요구사항 9.3).

본 파일은 빌드 순서상 가장 먼저 필요한 "저장 루트 결정 + templateId 검증 + 안전 경로
조립"만 구현한다(설계 §구성요소 1, 요구사항 2.1/2.3/2.4/2.7/2.8). 등록(register)·조회
(list/get)·삭제(delete)·Style_Profile 추출(extract)은 후속 작업에서 같은 모듈에 추가된다.

핵심 보안/격리 원칙:
- 모든 산출물은 `userData/templates/{templateId}/` 디렉토리 *내부에만* 기록한다(요구사항 2.1).
- 저장 루트는 `_resolve_local_root`(server.py)와 동일한 우선순위를 재사용하되, 템플릿은
  영속화가 필요하므로 tempdir로 폴백하지 않는다. AE_GENERATED_ROOT → ~/.agentic-editor →
  (둘 다 불가 시) None. None이면 호출자가 "no-storage-root" 에러를 반환한다(요구사항 2.4).
- templateId는 경로 조립 전에 항상 검증하고, realpath 기반으로 디렉토리 탈출을 차단한다
  (요구사항 2.7). server.py는 무겁고 import 부작용/순환 위험이 있으므로 import 하지 않고
  작은 우선순위 로직만 여기에 복제한다.
"""

from __future__ import annotations

import json
import os

# 템플릿 산출물이 저장되는 store_root 하위 서브디렉토리 이름.
# 최종 경로 형태: {store_root}/templates/{templateId}/{fname}
TEMPLATES_SUBDIR = "templates"

# templateId 길이 제약 (요구사항 2.7).
_TEMPLATE_ID_MIN_LEN = 1
_TEMPLATE_ID_MAX_LEN = 128


def resolve_template_store_root() -> str | None:
    """Template_Store 루트(=`templates/` 서브디렉토리가 놓일 userData 기반 디렉토리)를 반환.

    `_resolve_local_root`(server.py)의 우선순위를 재사용하되, 템플릿은 영속화되어야 하므로
    tempdir로 폴백하지 않는다(설계 §구성요소 1, 요구사항 2.3/2.4/2.8).

    우선순위:
    1. AE_GENERATED_ROOT 환경변수 (Electron이 주입하는 사용자별 userData 경로) —
       설정되어 있고 비어 있지 않으며 생성/쓰기 가능하면 사용.
    2. ~/.agentic-editor (사용자 홈, OS user별 격리, 항상 쓰기 가능 기대) — 생성 가능하면 사용.
    3. 둘 다 불가 → None. 호출자는 "no-storage-root" 에러를 반환한다(요구사항 2.4).

    Returns:
        store_root 절대 경로 문자열, 또는 결정 불가 시 None.
    """
    # 1) Electron 주입 userData 루트
    env_root = os.environ.get("AE_GENERATED_ROOT", "").strip()
    if env_root:
        try:
            os.makedirs(env_root, exist_ok=True)
            return env_root
        except (OSError, PermissionError):
            pass

    # 2) ~/.agentic-editor — OS user별 격리, 영속
    home_root = os.path.expanduser("~/.agentic-editor")
    try:
        os.makedirs(home_root, exist_ok=True)
        return home_root
    except (OSError, PermissionError):
        pass

    # 3) 결정 불가 — 템플릿은 tempdir로 폴백하지 않는다(영속 필요). 호출자가 no-storage-root 반환.
    return None


def _validate_template_id(tid: str) -> bool:
    """templateId가 경로 조립에 안전한지 검증한다(요구사항 2.7).

    유효 조건(AND):
    - 1자 이상 128자 이하
    - 경로 구분자 `/` 미포함
    - 윈도우 경로 구분자 `\\` 미포함
    - 상위 참조 `..` 미포함

    Returns:
        유효하면 True, 아니면 False.
    """
    if not isinstance(tid, str):
        return False
    if not (_TEMPLATE_ID_MIN_LEN <= len(tid) <= _TEMPLATE_ID_MAX_LEN):
        return False
    if "/" in tid or "\\" in tid or ".." in tid:
        return False
    return True


def safe_template_artifact_path(
    store_root: str, template_id: str, fname: str
) -> str | None:
    """`{store_root}/templates/{templateId}/{fname}` 안전 경로를 조립해 반환한다.

    templateId 검증을 먼저 수행하고(요구사항 2.7), realpath 기준으로 결과 경로가
    `{store_root}/templates/{templateId}` 디렉토리 *내부*에 머무는지 확인한다. 경로가
    해당 디렉토리를 벗어나면(예: fname에 `..`/절대경로 포함) None을 반환해 거부한다
    (요구사항 2.1).

    Args:
        store_root: `resolve_template_store_root()`가 반환한 userData 기반 루트.
        template_id: 검증 대상 templateId.
        fname: 디렉토리 내부 파일 이름(예: "base.pptx", "style_profile.json").

    Returns:
        안전한 절대/상대 경로 문자열, 또는 검증/탈출 차단 시 None.
    """
    if not store_root or not isinstance(store_root, str):
        return None
    if not _validate_template_id(template_id):
        return None
    if not isinstance(fname, str) or not fname:
        return None

    base_dir = os.path.join(store_root, TEMPLATES_SUBDIR, template_id)
    candidate = os.path.join(base_dir, fname)

    real_base = os.path.realpath(base_dir)
    real_candidate = os.path.realpath(candidate)

    # 결과가 base_dir 자신이거나 그 하위 경로일 때만 허용 — 그 외(탈출)는 거부.
    if real_candidate == real_base or real_candidate.startswith(real_base + os.sep):
        return candidate
    return None


# ===========================================================================
# Style_Profile 추출 (테마 XML) — 설계 §구성요소 2, 요구사항 3.1/3.2
# ===========================================================================
#
# 기준 `.pptx`의 테마 part XML에서 6개 토큰(주/보조/강조 색, 텍스트 색, 배경 색,
# 제목/본문 폰트)을 파싱한다. python-pptx는 테마를 공개 API로 직접 노출하지 않으므로
# 슬라이드 마스터의 part 관계를 통해 테마 part를 찾고, 그 lxml 엘리먼트를 직접 읽는다.
#
# 본 태스크(3.1)는 *추출 메커니즘*만 담당한다 — 즉 테마에서 읽을 수 있는 값을 그대로
# (색상은 normalize_color로 `#RRGGBB` 정규화 시도, 폰트는 typeface 문자열 그대로) 7키
# 딕셔너리에 담아 반환하며, 일부 토큰은 None/무효일 수 있다. per-token 폴백(부재·무효
# 토큰을 SLIDE_DESIGN 기본값으로 채우는 것)과 style_profile.json 저장은 후속 태스크
# (3.2)에서 같은 모듈에 추가된다.
#
# python-pptx(및 그 의존성 lxml)는 *지연 import* 한다. 본 모듈은 python-pptx 미설치
# 환경에서도 import 되어야 하므로(요구사항 9.3 대비) `pptx.*` 를 모듈 최상단에서
# import 하지 않는다. 또한 extract_style_profile()은 어떤 경우에도 예외를 던지지 않고
# 7키 딕셔너리를 반환한다(설계 §개요 폴백 격리 원칙).

# 슬라이드 마스터 → 테마 part 를 잇는 관계(relationship) URI.
_THEME_REL = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
)

# extract_style_profile()이 항상 반환하는 7개 키(설계 §데이터 모델 Style_Profile 스키마).
# style_profile.STYLE_PROFILE_KEY_ORDER와 동일하지만, python-pptx/style_profile 미설치
# 상황에서도 안전하게 동작하도록 여기서 명시적으로 정의한다.
_STYLE_PROFILE_KEYS = (
    "primaryColor",
    "secondaryColor",
    "accentColor",
    "textColor",
    "backgroundColor",
    "headingFont",
    "bodyFont",
)

# Style_Profile 색상 필드 → 테마 clrScheme 자식 태그(a: 네임스페이스) 매핑.
#   primaryColor   ← accent1
#   secondaryColor ← accent2
#   accentColor    ← accent3
#   textColor      ← dk1  (슬라이드 colorMap의 tx1 별칭)
#   backgroundColor← lt1  (슬라이드 colorMap의 bg1 별칭)
_COLOR_TAG_MAP = {
    "primaryColor": "a:accent1",
    "secondaryColor": "a:accent2",
    "accentColor": "a:accent3",
    "textColor": "a:dk1",
    "backgroundColor": "a:lt1",
}

# CSS/테마 폰트 스택에서 의미 있는 패밀리를 고를 때 건너뛰는 제네릭 토큰(소문자 비교).
# 폴백 시 첫 토큰('-apple-system')이 아니라 실제 패밀리(예: 'Apple SD Gothic Neo')를
# 선택하기 위함이다(설계 §구성요소 2 참고 노트).
_GENERIC_FONT_TOKENS = frozenset(
    {
        "-apple-system",
        "blinkmacsystemfont",
        "system-ui",
        "ui-sans-serif",
        "ui-serif",
        "ui-monospace",
        "ui-rounded",
        "sans-serif",
        "serif",
        "monospace",
        "cursive",
        "fantasy",
        "emoji",
        "math",
        "fangsong",
        "inherit",
        "initial",
        "unset",
    }
)


def _load_normalize_color():
    """`normalize_color`를 dual-path로 지연 import 한다(slide_templates.py와 동일 관례).

    repo 루트에서 실행 시 `ai_engine.style_profile`, ai_engine/ 내부에서 실행 시
    `style_profile` 로 import 된다. 둘 다 실패하면 None을 반환한다.

    Returns:
        normalize_color 함수, 또는 import 불가 시 None.
    """
    try:
        from ai_engine.style_profile import normalize_color

        return normalize_color
    except ImportError:  # pragma: no cover - exercised when run from ai_engine/
        try:
            from style_profile import normalize_color  # type: ignore

            return normalize_color
        except ImportError:
            return None


def _theme_element(prs):
    """python-pptx Presentation 객체에서 테마(a:theme) 루트 lxml 엘리먼트를 반환한다.

    첫 번째 슬라이드 마스터의 part 관계 중 테마 관계(_THEME_REL)를 따라가 테마 part 를
    찾는다. python-pptx 버전에 따라 테마 part 가 두 형태로 나타난다.
      - XmlPart 계열: `._element`(또는 `.element`)로 파싱된 lxml 루트를 직접 노출.
      - 일반 Part: lxml 루트가 없고 `.blob`(원본 XML 바이트)만 노출 → lxml로 파싱.
    전자를 우선 시도하고, 없으면 `.blob`을 lxml로 파싱한다. 마스터가 없거나 테마 관계를
    해석할 수 없으면(손상/비표준 .pptx 등) None을 반환하고 예외를 던지지 않는다.

    Args:
        prs: python-pptx Presentation 객체.

    Returns:
        a:theme 루트 lxml 엘리먼트, 또는 해석 불가 시 None.
    """
    try:
        master = prs.slide_masters[0]
        theme_part = master.part.part_related_by(_THEME_REL)
    except Exception:
        return None

    # 1) XmlPart 계열: 파싱된 lxml 루트를 직접 노출(_element 우선, element 보조).
    for attr in ("_element", "element"):
        try:
            elem = getattr(theme_part, attr, None)
            if elem is not None and hasattr(elem, "find"):
                return elem
        except Exception:
            pass

    # 2) 일반 Part: 원본 XML 바이트(blob)를 lxml로 직접 파싱.
    try:
        blob = getattr(theme_part, "blob", None)
        if blob:
            from lxml import etree

            return etree.fromstring(blob)
    except Exception:
        return None

    return None


def _color_from_clr_node(clr_node, qn):
    """clrScheme 자식 노드(예: <a:accent1>)에서 색상 16진수 문자열을 추출한다.

    OOXML에서 색은 두 형태로 나타난다.
      - <a:srgbClr val="RRGGBB"/>      → val 속성(6자리 16진수)
      - <a:sysClr val="windowText" lastClr="000000"/> → lastClr 속성(해석된 6자리)
    srgbClr를 우선 확인하고, 없으면 sysClr의 lastClr를 사용한다.

    Args:
        clr_node: clrScheme의 자식 색상 노드(없으면 None 가능).
        qn: pptx.oxml.ns.qn (네임스페이스 한정 태그 생성기).

    Returns:
        '#' 없는 6자리 16진수 문자열(예: '4472C4'), 또는 추출 불가 시 None.
    """
    if clr_node is None:
        return None
    srgb = clr_node.find(qn("a:srgbClr"))
    if srgb is not None:
        val = srgb.get("val")
        if val:
            return val
    sys_clr = clr_node.find(qn("a:sysClr"))
    if sys_clr is not None:
        last = sys_clr.get("lastClr")
        if last:
            return last
    return None


def _latin_typeface(font_node, qn):
    """majorFont/minorFont 노드에서 a:latin@typeface 문자열을 추출한다.

    Args:
        font_node: a:majorFont 또는 a:minorFont 노드(없으면 None 가능).
        qn: pptx.oxml.ns.qn.

    Returns:
        공백 제거한 typeface 문자열, 또는 부재/빈 문자열 시 None.
    """
    if font_node is None:
        return None
    latin = font_node.find(qn("a:latin"))
    if latin is None:
        return None
    typeface = latin.get("typeface")
    if isinstance(typeface, str) and typeface.strip():
        return typeface.strip()
    return None


def _first_real_family(stack_or_name):
    """폰트 스택 문자열에서 의미 있는 1–64자 단일 패밀리를 추출한다.

    SLIDE_DESIGN의 `font_heading`/`font_body`는 전체 CSS 폰트 스택
    (예: "-apple-system, BlinkMacSystemFont, 'Apple SD Gothic Neo', ...")이다.
    폴백 시 제네릭 선두 토큰('-apple-system' 등) 대신 실제 패밀리
    ('Apple SD Gothic Neo')를 고르기 위해, 콤마로 분리하고 따옴표를 제거한 뒤
    `_GENERIC_FONT_TOKENS`에 없는 첫 1–64자 패밀리를 반환한다. 실제 패밀리가 하나도
    없으면(=모두 제네릭) 1–64자 조건을 만족하는 첫 토큰을 마지막 수단으로 반환한다.

    Args:
        stack_or_name: 단일 패밀리명 또는 콤마 구분 폰트 스택 문자열.

    Returns:
        1–64자 단일 패밀리 문자열, 또는 추출 불가 시 None.
    """
    if not isinstance(stack_or_name, str):
        return None
    last_resort = None
    for part in stack_or_name.split(","):
        family = part.strip().strip('"').strip("'").strip()
        if not family or not (1 <= len(family) <= 64):
            continue
        if last_resort is None:
            last_resort = family
        if family.lower() in _GENERIC_FONT_TOKENS:
            continue
        return family
    return last_resort


def extract_style_profile(prs) -> dict:
    """python-pptx Presentation 객체의 테마에서 6개 스타일 토큰을 추출한다(요구사항 3.1, 3.2).

    추출 매핑(설계 §구성요소 2):
      - 색상(a:clrScheme):
          primaryColor    ← a:accent1
          secondaryColor  ← a:accent2
          accentColor     ← a:accent3
          textColor       ← a:dk1 (tx1)
          backgroundColor ← a:lt1 (bg1)
        각 색은 srgbClr@val 또는 sysClr@lastClr 에서 읽어 normalize_color로
        `#RRGGBB` 정규화를 시도한다(정규화 실패 시 해당 항목 None).
      - 폰트(a:fontScheme):
          headingFont ← a:majorFont/a:latin@typeface
          bodyFont    ← a:minorFont/a:latin@typeface

    반환 딕셔너리는 항상 7개 키(_STYLE_PROFILE_KEYS)를 가지며, 추출 불가하거나 형식이
    맞지 않는 토큰은 None 으로 남는다(부재·무효 토큰의 per-token 폴백은 후속 태스크 3.2
    에서 처리). 이 함수는 어떤 입력에도 예외를 던지지 않는다(설계 폴백 격리 원칙).

    Args:
        prs: python-pptx Presentation 객체.

    Returns:
        7개 키를 가진 dict. 각 값은 정규화된 색상/폰트 문자열 또는 None.
    """
    result = {key: None for key in _STYLE_PROFILE_KEYS}

    try:
        from pptx.oxml.ns import qn
    except ImportError:
        # python-pptx 미설치 — 그래도 7키 dict를 반환(요구사항 9.3 대비, 예외 금지).
        return result

    normalize_color = _load_normalize_color()

    theme = _theme_element(prs)
    if theme is None:
        return result

    # --- 색상 (a:clrScheme) ---
    try:
        clr_scheme = theme.find(".//" + qn("a:clrScheme"))
        if clr_scheme is not None:
            for field, tag in _COLOR_TAG_MAP.items():
                node = clr_scheme.find(qn(tag))
                raw = _color_from_clr_node(node, qn)
                if raw is None:
                    continue
                # normalize_color는 선택적 '#' + 6자리 16진수를 대문자 '#RRGGBB'로 변환.
                # 테마 값은 보통 '#' 없는 6자리이므로 그대로 정규화된다.
                if normalize_color is not None:
                    result[field] = normalize_color(raw)
                else:  # style_profile import 불가 — 최소한 형식만 맞으면 보존
                    result[field] = raw
    except Exception:
        # 색상 파싱 중 어떤 예외도 격리 — 폰트 추출은 계속 시도한다.
        pass

    # --- 폰트 (a:fontScheme) ---
    try:
        font_scheme = theme.find(".//" + qn("a:fontScheme"))
        if font_scheme is not None:
            major = font_scheme.find(qn("a:majorFont"))
            minor = font_scheme.find(qn("a:minorFont"))
            result["headingFont"] = _latin_typeface(major, qn)
            result["bodyFont"] = _latin_typeface(minor, qn)
    except Exception:
        pass

    return result


# ===========================================================================
# per-token 폴백 + style_profile.json 저장/조회 — 설계 §구성요소 2, 요구사항 3.3/3.4/3.5
# ===========================================================================
#
# extract_style_profile()는 테마에서 읽을 수 있는 값만 7키 dict에 담아 반환하며 일부
# 토큰은 None/무효일 수 있다(태스크 3.1). 여기서는 그 결과에 per-token 폴백을 적용해
# (부재·무효 토큰을 SLIDE_DESIGN 대응 기본값으로 채워) 7토큰이 항상 비어 있지 않은
# 완성된 Style_Profile dict를 만들고(build_style_profile), 이를 결정론적 JSON으로
# style_profile.json에 저장하고(save_style_profile), 매 호출 바이트 동일하게 다시
# 읽는다(get_style_profile, 요구사항 3.5).
#
# SLIDE_DESIGN / StyleProfile / serialize 는 slide_templates.py·style_profile.py에서
# dual-path 로 지연 import 한다(repo 루트 실행 시 ai_engine.*, ai_engine/ 내부 실행 시
# top-level). python-pptx 의존 코드는 extract_style_profile() 내부에서만 lazy import 된다.

# style_profile.json 파일명 — 디렉토리 내부 상대 이름(설계 §데이터 모델 디렉토리 레이아웃).
_STYLE_PROFILE_FILENAME = "style_profile.json"

# per-token 폴백 매핑 (요구사항 3.3): Style_Profile 색상 필드 → SLIDE_DESIGN 키.
_COLOR_FALLBACK_KEYS = {
    "primaryColor": "primary",
    "secondaryColor": "secondary",
    "accentColor": "accent",
    "textColor": "text_dark",
    "backgroundColor": "bg_light",
}

# Style_Profile 폰트 필드 → SLIDE_DESIGN 폰트 스택 키.
_FONT_FALLBACK_KEYS = {
    "headingFont": "font_heading",
    "bodyFont": "font_body",
}

# SLIDE_DESIGN 을 import 할 수 없는 환경에서 쓰는 로컬 폴백 값. slide_templates.SLIDE_DESIGN의
# 대응 값과 일치한다(색상은 #RRGGBB, 폰트는 _first_real_family로 단일 패밀리 추출 가능한 값).
_LOCAL_SLIDE_DESIGN_FALLBACK = {
    "primary": "#0066FF",
    "secondary": "#00C896",
    "accent": "#FF6B35",
    "text_dark": "#1A1A1A",
    "bg_light": "#FAFAFA",
    "font_heading": "Apple SD Gothic Neo",
    "font_body": "Apple SD Gothic Neo",
}


def _load_slide_design():
    """`SLIDE_DESIGN`을 dual-path로 지연 import 한다(slide_templates.py).

    repo 루트 실행 시 `ai_engine.slide_templates`, ai_engine/ 내부 실행 시 `slide_templates`.
    둘 다 실패하면 None 을 반환하고, 호출부는 `_LOCAL_SLIDE_DESIGN_FALLBACK`을 사용한다.

    Returns:
        SLIDE_DESIGN dict, 또는 import 불가 시 None.
    """
    try:
        from ai_engine.slide_templates import SLIDE_DESIGN

        return SLIDE_DESIGN
    except ImportError:  # pragma: no cover - exercised when run from ai_engine/
        try:
            from slide_templates import SLIDE_DESIGN  # type: ignore

            return SLIDE_DESIGN
        except ImportError:
            return None


def _load_style_profile_serializer():
    """`StyleProfile`과 `serialize`를 dual-path로 지연 import 한다(style_profile.py).

    Returns:
        (StyleProfile, serialize) 튜플, 또는 import 불가 시 (None, None).
    """
    try:
        from ai_engine.style_profile import StyleProfile, serialize

        return StyleProfile, serialize
    except ImportError:  # pragma: no cover - exercised when run from ai_engine/
        try:
            from style_profile import StyleProfile, serialize  # type: ignore

            return StyleProfile, serialize
        except ImportError:
            return None, None


def _fallback_normalize_color(value):
    """`normalize_color`를 import 할 수 없을 때 쓰는 최소 정규화(선택적 '#' + 6자리 16진수).

    Args:
        value: 정규화 대상 문자열.

    Returns:
        대문자 '#RRGGBB' 문자열, 또는 형식 불일치 시 None.
    """
    if not isinstance(value, str):
        return None
    digits = value[1:] if value.startswith("#") else value
    if len(digits) == 6 and all(c in "0123456789abcdefABCDEF" for c in digits):
        return "#" + digits.upper()
    return None


def build_style_profile(prs) -> dict:
    """테마에서 추출한 토큰에 per-token 폴백을 적용해 완성된 Style_Profile dict를 만든다.

    extract_style_profile(prs)로 7키(raw/None) dict를 얻은 뒤, 각 토큰을 개별 검증한다
    (요구사항 3.3):
      - 색상 5종(primary/secondary/accent/text/background): normalize_color로 `#RRGGBB`
        (대문자) 정규화에 성공하면 그 값을, 부재·무효면 SLIDE_DESIGN 대응 기본값
        (primary/secondary/accent/text_dark/bg_light)을 정규화해 채운다.
      - 폰트 2종(heading/body): _first_real_family로 1–64자 단일 패밀리를 얻으면 그 값을,
        부재·무효면 SLIDE_DESIGN 폰트 스택(font_heading/font_body)에서 _first_real_family로
        추출한 단일 패밀리를 채운다.
    결과 7토큰은 항상 비어 있지 않으며, 색상은 모두 정규화된 `#RRGGBB`, 폰트는 모두 1–64자
    문자열이다(요구사항 3.1 보장). 이 함수는 예외를 던지지 않는다(폴백 격리 원칙).

    Args:
        prs: python-pptx Presentation 객체.

    Returns:
        STYLE_PROFILE_KEY_ORDER와 동일한 7개 키를 가진 완성된 dict(모든 값 비어있지 않음).
    """
    raw = extract_style_profile(prs)

    normalize_color = _load_normalize_color()

    def _norm(value):
        if normalize_color is not None:
            return normalize_color(value) if isinstance(value, str) else None
        return _fallback_normalize_color(value)

    slide_design = _load_slide_design()
    if not isinstance(slide_design, dict):
        slide_design = _LOCAL_SLIDE_DESIGN_FALLBACK

    def _design_value(design_key):
        value = slide_design.get(design_key)
        if value is None:
            value = _LOCAL_SLIDE_DESIGN_FALLBACK.get(design_key)
        return value

    result: dict = {}

    # 색상 5종 — 추출값 정규화 시도 → 실패 시 SLIDE_DESIGN 기본값 정규화 → 로컬 상수 폴백.
    for field, design_key in _COLOR_FALLBACK_KEYS.items():
        normalized = _norm(raw.get(field))
        if normalized is None:
            normalized = _norm(_design_value(design_key))
        if normalized is None:
            # 최후 폴백 — 로컬 상수는 이미 #RRGGBB 형식.
            normalized = _fallback_normalize_color(
                _LOCAL_SLIDE_DESIGN_FALLBACK.get(design_key)
            )
        result[field] = normalized

    # 폰트 2종 — 추출 typeface에서 단일 패밀리 추출 → 실패 시 SLIDE_DESIGN 스택에서 추출.
    for field, design_key in _FONT_FALLBACK_KEYS.items():
        extracted = raw.get(field)
        family = _first_real_family(extracted) if isinstance(extracted, str) else None
        if family is None:
            family = _first_real_family(_design_value(design_key))
        if family is None:
            family = _first_real_family(
                _LOCAL_SLIDE_DESIGN_FALLBACK.get(design_key)
            )
        result[field] = family

    return result


def save_style_profile(profile_dict: dict, store_root: str, template_id: str) -> str:
    """완성된 Style_Profile dict를 결정론적 JSON으로 style_profile.json에 저장한다(요구사항 3.4).

    dict로부터 StyleProfile 객체를 만들고 ai_engine.style_profile.serialize()로 직렬화해
    `{store_root}/templates/{templateId}/style_profile.json`에 UTF-8 바이트로 쓴다. 경로는
    safe_template_artifact_path로 조립해 디렉토리 탈출을 차단한다(요구사항 2.1). serialize()가
    키 순서·공백을 고정하므로 동일 입력 → 매 저장마다 바이트 동일하다(요구사항 4.2 → 3.5 기반).

    Args:
        profile_dict: build_style_profile()가 반환한 7키 완성 dict.
        store_root: resolve_template_store_root()가 반환한 userData 기반 루트.
        template_id: 대상 templateId(검증 대상).

    Returns:
        저장된 style_profile.json의 경로 문자열.

    Raises:
        RuntimeError: style_profile 모듈을 import 할 수 없는 경우.
        ValueError: templateId 검증 실패 등으로 안전 경로를 조립할 수 없는 경우.
        OSError: 디렉토리 생성/파일 쓰기 실패(호출부 register가 template-store-write-failed로 변환).
    """
    StyleProfile, serialize = _load_style_profile_serializer()
    if StyleProfile is None or serialize is None:
        raise RuntimeError("style_profile module unavailable")

    path = safe_template_artifact_path(store_root, template_id, _STYLE_PROFILE_FILENAME)
    if path is None:
        raise ValueError("invalid-template-id")

    profile = StyleProfile(
        primaryColor=profile_dict["primaryColor"],
        secondaryColor=profile_dict["secondaryColor"],
        accentColor=profile_dict["accentColor"],
        textColor=profile_dict["textColor"],
        backgroundColor=profile_dict["backgroundColor"],
        headingFont=profile_dict["headingFont"],
        bodyFont=profile_dict["bodyFont"],
    )
    text = serialize(profile)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    with open(path, "wb") as fh:
        fh.write(text.encode("utf-8"))
    return path


def get_style_profile(template_id: str, store_root: str) -> dict:
    """저장된 style_profile.json을 읽어 파싱된 dict를 반환한다(요구사항 3.5).

    style_profile.json은 serialize()로 결정론적으로 기록되므로, 저장된 바이트를 읽어
    파싱해 돌려주는 동작은 매 호출 바이트 단위로 안정적이다(동일 templateId → 동일 내용).

    Args:
        template_id: 조회 대상 templateId(검증 대상).
        store_root: resolve_template_store_root()가 반환한 userData 기반 루트.

    Returns:
        파싱된 Style_Profile dict(7키), 또는 실패 시 {"error": ...}.
        error ∈ {invalid-template-id, template-not-found, invalid-json}.
    """
    if not _validate_template_id(template_id):
        return {"error": "invalid-template-id"}

    path = safe_template_artifact_path(store_root, template_id, _STYLE_PROFILE_FILENAME)
    if path is None:
        return {"error": "invalid-template-id"}

    if not os.path.isfile(path):
        return {"error": "template-not-found"}

    try:
        with open(path, "rb") as fh:
            raw_bytes = fh.read()
        return json.loads(raw_bytes.decode("utf-8"))
    except (OSError, ValueError, UnicodeDecodeError):
        # JSONDecodeError 는 ValueError 의 서브타입이므로 함께 포착된다.
        return {"error": "invalid-json"}


# ===========================================================================
# 템플릿 등록 (register) — 설계 §구성요소 1(등록 검증 순서), 요구사항 1, 2.2/2.5/2.6, 9.3
# ===========================================================================
#
# register_template()은 사용자가 제출한 `.pptx` 파일을 Template_Store에 안전하게
# 등록한다. 검증은 *짧은 작업 먼저, 디스크 쓰기는 마지막에* 수행해 부분 산출물을
# 방지한다(설계 §구성요소 1). 이 함수는 어떤 입력/오류에도 예외를 밖으로 던지지 않고
# 항상 dict를 반환한다(폴백 격리 원칙).
#
# 검증 순서(요구사항 매핑):
#   1. python-pptx import 불가          → missing-dep                (9.3)
#   2. 이름 trim 길이 ∉ [1,100]         → invalid-name               (1.2, 1.7)
#   3. 파일 크기 > 50MB                 → template-too-large         (1.4)
#      (파일 부재/접근 불가 시 invalid-template로 처리)
#   4. 확장자 .pptx AND Presentation 열림 실패 → invalid-template     (1.3)
#   5. store_root 결정 불가             → no-storage-root            (2.4)
#   6. 이름 중복(trim + casefold)       → duplicate-name             (1.6)
#   7. {templateId}/ 생성 + base.pptx 복사 + style_profile.json + metadata.json
#      중간 예외 시 부분 산출물 정리 후  → template-store-write-failed (2.6)
#   8. 성공                              → {templateId, name, path, layoutCount} (1.5)

# ===========================================================================
# 템플릿 구조 분석 (레이아웃 역할·placeholder) — 모델이 콘텐츠를 적합한 레이아웃에
# 매핑하도록 프롬프트에 제공할 구조 정보를 추출한다. server.py의 _layout_name_matches /
# _layout_has_content_placeholder와 동일한 역할 어휘를 (import 없이) 복제한다.
# ===========================================================================

# 레이아웃 역할 분류 — server.py와 동일 의미. 'two-column'이 'content'보다 먼저 검사된다
# (비교/2단이 'content' 키워드도 포함할 수 있으므로 더 구체적인 역할을 우선).
_LAYOUT_ROLE_KEYWORDS = (
    ("title", ("표지", "cover", "title slide", "제목 슬라이드")),
    ("two-column", ("two content", "two-content", "twocontent", "two column",
                    "two-column", "comparison", "비교", "2단")),
    ("section", ("section header", "section", "섹션", "구분", "divider", "간지")),
    ("content", ("title and content", "제목 및 내용", "content", "내용")),
    ("blank", ("blank", "빈", "공백")),
)


def _classify_layout_role(name_norm: str, has_content: bool) -> str:
    """정규화된 레이아웃 이름 + 콘텐츠 placeholder 보유 여부로 역할을 분류한다.

    server.py의 _layout_name_matches 어휘와 일관된다. 'title'은 '내용/content'를 포함하면
    표지로 보지 않는다(제목+내용 레이아웃 오인 방지). 매칭되는 키워드가 없으면 콘텐츠
    placeholder 보유 시 'content', 아니면 'other'로 분류한다.
    """
    n = name_norm or ""
    # title: '내용/content' 포함 시 표지 아님
    if ("표지" in n or "cover" in n or "title slide" in n or "제목 슬라이드" in n) \
            and "content" not in n and "내용" not in n:
        return "title"
    for role, kws in _LAYOUT_ROLE_KEYWORDS:
        if role == "title":
            continue  # 위에서 별도 처리
        if any(kw in n for kw in kws):
            return role
    return "content" if has_content else "other"


def _placeholder_kinds(layout) -> list:
    """레이아웃의 placeholder 종류 이름 목록을 반환한다(예: ['TITLE', 'BODY', 'PICTURE']).

    python-pptx enum 사용 불가/예외 시 빈 목록을 반환한다(폴백 격리).
    """
    kinds = []
    try:
        for ph in layout.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf is None:
                continue
            try:
                t = pf.type
                name = getattr(t, "name", None) or str(t)
            except Exception:
                name = "UNKNOWN"
            kinds.append(name)
    except Exception:
        return []
    return kinds


def _layout_has_content(layout) -> bool:
    """레이아웃이 본문(body/object) placeholder(idx>=1)를 보유하는지 — server.py와 동일 기준."""
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        body_types = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
    except Exception:
        body_types = None
    try:
        for ph in layout.placeholders:
            pf = getattr(ph, "placeholder_format", None)
            if pf is None:
                continue
            idx = pf.idx
            if idx is None or idx < 1:
                continue
            if body_types is None:
                return True
            try:
                if pf.type in body_types:
                    return True
            except Exception:
                continue
    except Exception:
        return False
    return False


def extract_template_structure(prs) -> dict:
    """템플릿의 슬라이드 레이아웃 구조를 추출한다(모델 프롬프트용).

    각 레이아웃에 대해 {index, name, role, placeholders, hasContent}를 수집한다. role은
    표지/2단/비교(two-column)/섹션(section)/내용(content)/빈(blank)/기타(other)로 분류된다.
    역할별 개수 요약(roleCounts)도 함께 반환해 모델이 "이 템플릿엔 표지 1, 2단 2, 섹션 3개가
    있다"처럼 콘텐츠를 매핑할 수 있게 한다.

    어떤 예외에도 raise하지 않고 {"layouts": [], "roleCounts": {}}로 폴백한다(격리).

    Returns:
        {"layouts": [{index, name, role, placeholders, hasContent}, ...],
         "roleCounts": {role: count, ...}}
    """
    result = {"layouts": [], "roleCounts": {}}
    try:
        layouts = list(prs.slide_layouts)
    except Exception:
        return result

    for idx, layout in enumerate(layouts):
        try:
            name = (getattr(layout, "name", "") or "").strip()
        except Exception:
            name = ""
        name_norm = name.lower()
        has_content = _layout_has_content(layout)
        role = _classify_layout_role(name_norm, has_content)
        kinds = _placeholder_kinds(layout)
        result["layouts"].append({
            "index": idx,
            "name": name or f"Layout {idx}",
            "role": role,
            "placeholders": kinds,
            "hasContent": has_content,
        })
        result["roleCounts"][role] = result["roleCounts"].get(role, 0) + 1

    return result


# 기준 .pptx 복사본 파일명(디렉토리 내부 상대 이름).
_BASE_PPTX_FILENAME = "base.pptx"

# 등록 메타데이터 파일명(디렉토리 내부 상대 이름).
_METADATA_FILENAME = "metadata.json"

# 템플릿 이름 trim 후 허용 길이 범위(요구사항 1.2, 1.7).
_NAME_MIN_LEN = 1
_NAME_MAX_LEN = 100

# 업로드 허용 최대 파일 크기(바이트). 50MB(요구사항 1.4).
_MAX_TEMPLATE_BYTES = 52428800

# 원인 설명 문자열 최대 길이(요구사항 1.3/2.6 — "≤200자").
_DETAIL_MAX_LEN = 200


def _iter_template_metadata(store_root):
    """`{store_root}/templates/*/metadata.json`을 순회하며 파싱된 메타데이터 dict를 yield한다.

    register_template의 이름 중복 검사(요구사항 1.6)에 쓰이며, list_templates(태스크 4.1)도
    재사용할 수 있도록 제너레이터로 둔다. 읽을 수 없거나 JSON 파싱에 실패하거나 dict가 아닌
    항목은 조용히 건너뛴다(손상 항목이 등록/목록 전체를 막지 않도록).

    Args:
        store_root: resolve_template_store_root()가 반환한 userData 기반 루트.

    Yields:
        각 템플릿 디렉토리의 metadata.json을 파싱한 dict.
    """
    if not store_root or not isinstance(store_root, str):
        return
    templates_dir = os.path.join(store_root, TEMPLATES_SUBDIR)
    try:
        entries = os.listdir(templates_dir)
    except OSError:
        return
    for entry in entries:
        meta_path = os.path.join(templates_dir, entry, _METADATA_FILENAME)
        if not os.path.isfile(meta_path):
            continue
        try:
            with open(meta_path, "rb") as fh:
                data = json.loads(fh.read().decode("utf-8"))
        except (OSError, ValueError, UnicodeDecodeError):
            # 손상/읽기 불가 항목은 건너뛴다(JSONDecodeError는 ValueError 서브타입).
            continue
        if isinstance(data, dict):
            yield data


def register_template(file_path: str, name: str, store_root: str = None) -> dict:
    """`.pptx` 파일을 검증 후 Template_Store에 등록한다(요구사항 1, 2.2/2.5/2.6, 9.3).

    검증 순서는 모듈 상단 주석 및 설계 §구성요소 1의 "등록 검증 순서"를 따른다. 짧은 검사
    (의존성·이름·크기·열기·중복)를 먼저 수행하고, 디스크 쓰기(디렉토리 생성·파일 복사·
    Style_Profile/metadata 저장)는 가장 마지막에 한 번에 수행한다. 쓰기 도중 어떤 예외가
    발생하면 부분적으로 생성된 `{templateId}` 디렉토리를 정리하고 template-store-write-failed를
    반환한다(요구사항 2.6). 이 함수는 어떤 경우에도 예외를 밖으로 던지지 않는다.

    Args:
        file_path: 등록할 `.pptx` 파일의 (로컬) 절대/상대 경로.
        name: 사용자가 입력한 템플릿 이름(앞뒤 공백 제거 후 1–100자여야 함).
        store_root: 저장 루트. None이면 resolve_template_store_root()로 결정한다.

    Returns:
        성공: {"templateId", "name", "path"(상대), "layoutCount"} (요구사항 1.5)
        실패: {"error", ...} —
            missing-dep / invalid-name / template-too-large / invalid-template /
            no-storage-root / duplicate-name / template-store-write-failed
    """
    try:
        # --- 1) python-pptx import 가능 여부 (요구사항 9.3) ---
        try:
            from pptx import Presentation
        except ImportError:
            return {
                "error": "missing-dep",
                "lib": "python-pptx",
                "hint": "pip install python-pptx",
            }

        # --- 2) 이름 trim 후 길이 1–100 (요구사항 1.2, 1.7) ---
        if not isinstance(name, str):
            return {"error": "invalid-name", "allowed": [_NAME_MIN_LEN, _NAME_MAX_LEN]}
        name = name.strip()
        if not (_NAME_MIN_LEN <= len(name) <= _NAME_MAX_LEN):
            return {"error": "invalid-name", "allowed": [_NAME_MIN_LEN, _NAME_MAX_LEN]}

        # --- 3) 파일 크기 ≤ 50MB (요구사항 1.4). 부재/접근 불가는 invalid-template로 처리. ---
        if not isinstance(file_path, str) or not file_path:
            return {"error": "invalid-template", "detail": "파일 경로가 유효하지 않습니다"}
        try:
            size = os.path.getsize(file_path)
        except OSError as exc:
            return {"error": "invalid-template", "detail": str(exc)[:_DETAIL_MAX_LEN]}
        if size > _MAX_TEMPLATE_BYTES:
            return {"error": "template-too-large", "maxBytes": _MAX_TEMPLATE_BYTES}

        # --- 4) 확장자 .pptx AND Presentation()로 열림 (요구사항 1.3). 실패 시 저장 안 함. ---
        if os.path.splitext(file_path)[1].lower() != ".pptx":
            return {
                "error": "invalid-template",
                "detail": "확장자가 .pptx가 아닙니다",
            }
        try:
            prs = Presentation(file_path)
        except Exception as exc:
            return {"error": "invalid-template", "detail": str(exc)[:_DETAIL_MAX_LEN]}

        # --- 5) store_root 결정 (요구사항 2.4). ---
        if not store_root:
            store_root = resolve_template_store_root()
        if not store_root:
            return {"error": "no-storage-root"}

        # --- 6) 이름 중복 검사: trim + casefold 비교 (요구사항 1.6). 기존 산출물 보존. ---
        target_key = name.casefold()
        for meta in _iter_template_metadata(store_root):
            existing = meta.get("name")
            if isinstance(existing, str) and existing.strip().casefold() == target_key:
                return {"error": "duplicate-name", "name": name}

        # --- 7) 디스크 쓰기 (요구사항 2.2, 2.5, 2.6). 중간 예외 시 부분 산출물 정리. ---
        import shutil
        import uuid
        from datetime import datetime, timezone

        template_id = str(uuid.uuid4())
        template_dir = os.path.join(store_root, TEMPLATES_SUBDIR, template_id)
        try:
            base_path = safe_template_artifact_path(
                store_root, template_id, _BASE_PPTX_FILENAME
            )
            meta_path = safe_template_artifact_path(
                store_root, template_id, _METADATA_FILENAME
            )
            if base_path is None or meta_path is None:
                # uuid4 → 정상적으로는 도달 불가. 방어적으로 write-failed 처리.
                raise ValueError("invalid-template-id")

            os.makedirs(template_dir, exist_ok=True)

            # base.pptx 복사 — 원본 경로를 신뢰해 그 밖으로 쓰지 않고 store_root 하위로만 복사.
            shutil.copyfile(file_path, base_path)

            # Style_Profile 추출 + 결정론적 JSON 저장 (요구사항 3.x — 기존 헬퍼 재사용).
            profile_dict = build_style_profile(prs)
            save_style_profile(profile_dict, store_root, template_id)

            # 레이아웃 구조 추출 — 모델이 콘텐츠를 적합한 레이아웃에 매핑하도록 metadata에 저장.
            try:
                structure = extract_template_structure(prs)
            except Exception:
                structure = {"layouts": [], "roleCounts": {}}

            layout_count = len(prs.slide_layouts)
            created_at = datetime.now(timezone.utc).isoformat().replace("+00:00", "Z")
            metadata = {
                "templateId": template_id,
                "name": name,
                "createdAt": created_at,
                "basePptx": _BASE_PPTX_FILENAME,
                "layoutCount": layout_count,
                "layouts": structure.get("layouts", []),
                "roleCounts": structure.get("roleCounts", {}),
            }
            with open(meta_path, "wb") as fh:
                fh.write(json.dumps(metadata, ensure_ascii=False).encode("utf-8"))
        except Exception as exc:
            # 부분적으로 생성된 디렉토리를 남기지 않는다(요구사항 2.6).
            shutil.rmtree(template_dir, ignore_errors=True)
            return {
                "error": "template-store-write-failed",
                "detail": str(exc)[:_DETAIL_MAX_LEN],
            }

        # --- 8) 성공 (요구사항 1.5). path는 store_root 기준 상대 경로. ---
        return {
            "templateId": template_id,
            "name": name,
            "path": "{}/{}/{}".format(
                TEMPLATES_SUBDIR, template_id, _BASE_PPTX_FILENAME
            ),
            "layoutCount": layout_count,
        }
    except Exception as exc:  # 최후 방어 — register_template은 절대 raise하지 않는다.
        return {
            "error": "template-store-write-failed",
            "detail": str(exc)[:_DETAIL_MAX_LEN],
        }


# ===========================================================================
# 템플릿 조회·삭제 (list / get / delete) — 설계 §구성요소 1, 요구사항 2.7/5.3/5.4/8.1/8.8/8.12
# ===========================================================================
#
# register_template()이 기록한 metadata.json({templateId, name, createdAt, basePptx,
# layoutCount})과 style_profile.json을 읽어 목록·단건 조회를 제공하고, 템플릿 디렉토리를
# 통째로 삭제한다. 세 함수 모두 어떤 입력/오류에도 예외를 밖으로 던지지 않고 항상 dict를
# 반환한다(폴백 격리 원칙). 모든 경로 조립은 _validate_template_id + realpath 탈출 차단을
# 거친다(요구사항 2.1/2.7).

# 목록 표시 상한(요구사항 8.1) — 최대 200개.
_LIST_MAX = 200


def list_templates(store_root: str = None) -> dict:
    """등록된 템플릿 목록을 등록 시각(createdAt) 내림차순으로 반환한다(요구사항 8.1).

    store_root가 None이면 resolve_template_store_root()로 결정한다. 저장 루트를 결정할 수
    없으면(아직 어떤 템플릿도 등록되지 않은 환경 등) 이는 목록 조회에서 오류가 아니므로
    빈 목록을 반환한다. _iter_template_metadata로 각 템플릿의 metadata.json을 순회하고,
    각 항목에서 {templateId, name, createdAt} 세 필드만 추려 createdAt 기준 내림차순으로
    정렬한 뒤 최대 200개(_LIST_MAX)로 자른다. createdAt은 ISO 8601 문자열이므로 사전식
    내림차순 정렬이 곧 최신순 정렬이다.

    Args:
        store_root: 저장 루트. None이면 resolve_template_store_root()로 결정한다.

    Returns:
        {"templates": [{"templateId", "name", "createdAt"}, ...]} — createdAt 내림차순,
        최대 200개. 저장 루트 부재 시 {"templates": []}.
    """
    if not store_root:
        store_root = resolve_template_store_root()
    if not store_root:
        # 저장 루트 없음 = 등록된 템플릿 없음. 목록 조회에서는 오류가 아니다.
        return {"templates": []}

    items = []
    for meta in _iter_template_metadata(store_root):
        tid = meta.get("templateId")
        if not isinstance(tid, str) or not tid:
            # templateId가 없는 손상 메타데이터는 목록에서 제외한다.
            continue
        items.append(
            {
                "templateId": tid,
                "name": meta.get("name"),
                "createdAt": meta.get("createdAt"),
            }
        )

    # createdAt(ISO 8601) 사전식 내림차순 = 최신 등록 순. None/누락은 빈 문자열로 취급.
    items.sort(key=lambda m: m.get("createdAt") or "", reverse=True)
    return {"templates": items[:_LIST_MAX]}


def get_template(template_id: str, store_root: str = None) -> dict:
    """단일 템플릿의 메타데이터·기준 경로·Style_Profile을 반환한다(요구사항 5.3).

    검증 순서:
      1. templateId 검증 실패 → {"error": "invalid-template-id"} (요구사항 2.7)
      2. store_root 결정(None이면 resolve). 결정 불가 또는 metadata.json 부재/손상 →
         {"error": "template-not-found"} (요구사항 5.4)
      3. 정상 → {templateId, name, templatePath(절대), styleProfile, createdAt}

    templatePath는 safe_template_artifact_path로 조립한 `{store_root}/templates/{id}/base.pptx`
    의 절대 경로다(디렉토리 탈출 차단). styleProfile은 get_style_profile()이 반환하는 dict로,
    style_profile.json 바이트를 매 호출 동일하게 파싱한 7키 dict다(요구사항 3.5).

    Args:
        template_id: 조회 대상 templateId(검증 대상).
        store_root: 저장 루트. None이면 resolve_template_store_root()로 결정한다.

    Returns:
        성공: {"templateId", "name", "templatePath", "styleProfile", "createdAt"}
        실패: {"error": "invalid-template-id" | "template-not-found"}
    """
    if not _validate_template_id(template_id):
        return {"error": "invalid-template-id"}

    if not store_root:
        store_root = resolve_template_store_root()
    if not store_root:
        # 저장 루트가 없으면 해당 템플릿 디렉토리도 존재할 수 없다.
        return {"error": "template-not-found"}

    meta_path = safe_template_artifact_path(store_root, template_id, _METADATA_FILENAME)
    base_path = safe_template_artifact_path(store_root, template_id, _BASE_PPTX_FILENAME)
    if meta_path is None or base_path is None:
        # _validate_template_id 통과 후에는 정상적으로 도달 불가. 방어적 처리.
        return {"error": "invalid-template-id"}

    if not os.path.isfile(meta_path):
        return {"error": "template-not-found"}

    try:
        with open(meta_path, "rb") as fh:
            metadata = json.loads(fh.read().decode("utf-8"))
    except (OSError, ValueError, UnicodeDecodeError):
        # 읽기 불가/손상 metadata.json → 조회 대상으로 취급하지 않는다.
        return {"error": "template-not-found"}
    if not isinstance(metadata, dict):
        return {"error": "template-not-found"}

    style_profile = get_style_profile(template_id, store_root)

    return {
        "templateId": template_id,
        "name": metadata.get("name"),
        "templatePath": os.path.abspath(base_path),
        "styleProfile": style_profile,
        "createdAt": metadata.get("createdAt"),
        "layoutCount": metadata.get("layoutCount"),
        "layouts": metadata.get("layouts", []),
        "roleCounts": metadata.get("roleCounts", {}),
    }


def delete_template(template_id: str, store_root: str = None) -> dict:
    """템플릿 디렉토리(`{store_root}/templates/{id}`)와 하위 전부를 제거한다(요구사항 8.8).

    검증 순서:
      1. templateId 검증 실패 → {"error": "invalid-template-id"} (요구사항 2.7)
      2. store_root 결정(None이면 resolve). 결정 불가 또는 디렉토리 부재 →
         {"error": "template-not-found"}
      3. shutil.rmtree로 제거. 성공 → {"ok": True, "templateId": id}
      4. 제거 중 예외 → 디렉토리를 보존한 채(ignore_errors 미사용)
         {"error": "template-delete-failed", "detail": "<원인 ≤200자>"} (요구사항 8.12)

    디렉토리 경로는 _validate_template_id 통과 후 realpath 기준으로 `{store_root}/templates`
    내부에 머무는지 재확인해 디렉토리 탈출을 차단한다(요구사항 2.1/2.7).

    Args:
        template_id: 삭제 대상 templateId(검증 대상).
        store_root: 저장 루트. None이면 resolve_template_store_root()로 결정한다.

    Returns:
        성공: {"ok": True, "templateId": template_id}
        실패: {"error": "invalid-template-id" | "template-not-found" | "template-delete-failed"}
    """
    if not _validate_template_id(template_id):
        return {"error": "invalid-template-id"}

    if not store_root:
        store_root = resolve_template_store_root()
    if not store_root:
        return {"error": "template-not-found"}

    templates_root = os.path.join(store_root, TEMPLATES_SUBDIR)
    template_dir = os.path.join(templates_root, template_id)

    # realpath 기준으로 결과 디렉토리가 templates/ 루트 내부에 머무는지 재확인(탈출 차단).
    real_root = os.path.realpath(templates_root)
    real_dir = os.path.realpath(template_dir)
    if real_dir != real_root and not real_dir.startswith(real_root + os.sep):
        return {"error": "invalid-template-id"}

    if not os.path.isdir(template_dir):
        return {"error": "template-not-found"}

    # 지연 import — 모듈 최상단 import 부작용을 피한다(파일 전반의 관례).
    import shutil

    try:
        # ignore_errors=False — 실패 시 예외를 받아 디렉토리를 보존한 채 에러를 반환한다.
        shutil.rmtree(template_dir)
    except Exception as exc:  # noqa: BLE001 — 모든 제거 실패를 단일 에러로 변환(요구사항 8.12)
        return {
            "error": "template-delete-failed",
            "detail": str(exc)[:_DETAIL_MAX_LEN],
        }

    return {"ok": True, "templateId": template_id}
