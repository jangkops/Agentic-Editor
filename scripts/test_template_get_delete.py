"""Task 4.2 검증 — Template_Manager 조회·삭제 단위 테스트 (요구사항 5.4 / 8.12).

직접 실행형: ai_engine/.venv/bin/python scripts/test_template_get_delete.py

검증 대상 모듈: ai_engine/template_manager.py (이미 구현 완료된
register_template / get_template / delete_template 사용 — 구현은 수정하지 않는다).

검증 항목:
  1. 부재 조회 (요구사항 5.4):
     등록되지 않은 templateId로 get_template 호출 → {"error": "template-not-found"}.
  2. 정상 등록 → 조회 → 삭제 사이클:
     유효 .pptx로 register_template → templateId 획득 → get_template이
     {templateId, name, templatePath, styleProfile, createdAt} 반환 →
     delete_template 성공 시 {ok, templateId} 반환 + templates/{id}/ 디렉토리와
     하위 파일이 실제로 제거됨.
  3. 삭제 예외 주입 시 디렉토리 보존 (요구사항 8.12):
     shutil.rmtree를 monkeypatch하여 예외를 던지게 하면 delete_template이
     디렉토리를 보존한 채 {"error": "template-delete-failed"}를 반환.

python-pptx import 불가 시 전체 테스트를 skip 한다(fixture 생성 불가).
"""
from __future__ import annotations

import os
import sys
import shutil
import tempfile
from unittest.mock import patch

# repo 루트를 import 경로에 추가 — `from ai_engine import template_manager` 가능하게 한다.
# (template_manager 내부 헬퍼는 ai_engine.* dual-path import 를 우선 시도한다.)
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

# fixture(.pptx) 생성에 python-pptx 가 필요하다. 미설치면 전체 테스트 skip.
try:
    from pptx import Presentation  # noqa: E402

    _HAS_PPTX = True
except ImportError:  # pragma: no cover - 설치 환경에서만 실행
    _HAS_PPTX = False

from ai_engine import template_manager as tm  # noqa: E402


# ---------- fixtures / helpers ----------

def _make_valid_pptx(path: str) -> None:
    """python-pptx로 최소 유효 .pptx를 생성한다(슬라이드 마스터/레이아웃/테마 포함)."""
    prs = Presentation()
    # 표지 레이아웃 1장 추가 — register_template의 Presentation 열기/레이아웃 카운트에 충분.
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)


def _template_dir(store_root: str, template_id: str) -> str:
    """{store_root}/templates/{templateId} 경로를 조립한다(검증 편의용)."""
    return os.path.join(store_root, tm.TEMPLATES_SUBDIR, template_id)


def _register_fixture(store_root: str, name: str = "My Template") -> str:
    """임시 store_root에 유효 .pptx를 등록하고 templateId를 반환한다."""
    src = os.path.join(store_root, "fixture.pptx")
    _make_valid_pptx(src)
    res = tm.register_template(src, name, store_root=store_root)
    assert "error" not in res, f"fixture 등록 실패: {res}"
    tid = res.get("templateId")
    assert isinstance(tid, str) and tid, f"templateId 누락: {res}"
    return tid


# ---------- tests ----------

def test_get_missing_returns_not_found():
    """부재 templateId 조회 → template-not-found (요구사항 5.4)."""
    store_root = tempfile.mkdtemp(prefix="ae-tpl-get-")
    try:
        # 어떤 템플릿도 등록되지 않은 store_root. 유효 형식이지만 미등록 id.
        res = tm.get_template("00000000-0000-4000-8000-000000000000", store_root=store_root)
        assert res == {"error": "template-not-found"}, res
        print("  [OK] 부재 templateId 조회 → template-not-found (요구사항 5.4)")
    finally:
        shutil.rmtree(store_root, ignore_errors=True)


def test_register_get_delete_cycle():
    """등록 → 조회 → 삭제 정상 사이클: 조회 필드 + 디렉토리 제거 검증."""
    store_root = tempfile.mkdtemp(prefix="ae-tpl-cycle-")
    try:
        tid = _register_fixture(store_root, name="Quarterly Deck")
        tdir = _template_dir(store_root, tid)
        assert os.path.isdir(tdir), f"등록 후 디렉토리 부재: {tdir}"

        # --- 조회: 5개 필드를 모두 반환 ---
        got = tm.get_template(tid, store_root=store_root)
        assert "error" not in got, got
        for key in ("templateId", "name", "templatePath", "styleProfile", "createdAt"):
            assert key in got, f"get_template 응답에 '{key}' 누락: {got}"
        assert got["templateId"] == tid, got
        assert got["name"] == "Quarterly Deck", got
        assert os.path.isabs(got["templatePath"]), f"templatePath 비절대경로: {got['templatePath']}"
        assert os.path.isfile(got["templatePath"]), f"base.pptx 부재: {got['templatePath']}"
        assert isinstance(got["styleProfile"], dict) and got["styleProfile"], got
        print("  [OK] get_template → templateId/name/templatePath/styleProfile/createdAt 반환")

        # --- 삭제: {ok, templateId} 반환 + 디렉토리/하위 제거 ---
        res = tm.delete_template(tid, store_root=store_root)
        assert res == {"ok": True, "templateId": tid}, res
        assert not os.path.exists(tdir), f"삭제 후 디렉토리 잔존: {tdir}"
        print("  [OK] delete_template → {ok, templateId} + 디렉토리/하위 제거 검증")
    finally:
        shutil.rmtree(store_root, ignore_errors=True)


def test_delete_exception_preserves_directory():
    """삭제 중 예외 주입 → 디렉토리 보존 + template-delete-failed (요구사항 8.12)."""
    store_root = tempfile.mkdtemp(prefix="ae-tpl-delfail-")
    try:
        tid = _register_fixture(store_root, name="Keep On Failure")
        tdir = _template_dir(store_root, tid)
        assert os.path.isdir(tdir), f"등록 후 디렉토리 부재: {tdir}"

        # delete_template 내부의 `import shutil; shutil.rmtree(...)` 가 예외를 만나도록 주입.
        with patch("shutil.rmtree", side_effect=OSError("injected delete failure")):
            res = tm.delete_template(tid, store_root=store_root)

        assert res.get("error") == "template-delete-failed", res
        # 예외 발생 시 디렉토리와 하위 산출물이 그대로 보존되어야 한다.
        assert os.path.isdir(tdir), f"삭제 실패 후 디렉토리 소실(보존 위반): {tdir}"
        assert os.path.isfile(os.path.join(tdir, "base.pptx")), "base.pptx 보존 위반"
        assert os.path.isfile(os.path.join(tdir, "metadata.json")), "metadata.json 보존 위반"
        print("  [OK] 삭제 예외 주입 → 디렉토리 보존 + template-delete-failed (요구사항 8.12)")
    finally:
        shutil.rmtree(store_root, ignore_errors=True)


def main():
    print("=== Task 4.2: Template_Manager 조회·삭제 단위 테스트 ===")
    if not _HAS_PPTX:
        print("  [SKIP] python-pptx 미설치 — fixture(.pptx) 생성 불가로 전체 skip")
        return
    test_get_missing_returns_not_found()
    test_register_get_delete_cycle()
    test_delete_exception_preserves_directory()
    print("ALL PASSED")


if __name__ == "__main__":
    main()
