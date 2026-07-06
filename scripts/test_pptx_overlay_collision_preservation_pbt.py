"""Preservation property tests — spec: pptx-overlay-collision-fix (bugfix), Task 2.

PROPERTY 2 — Preservation (비버그 입력의 배치 좌표/렌더 플랜 보존 — no-op 동등성).

이 PBT는 오버레이/충돌 수정(tasks 3.x)이 **변경해서는 안 되는** 동작을 고정한다.
**observation-first** 방법론을 따른다 — 각 단언은 먼저 *미수정* 코드(``ai_engine``)가
비버그(``isBugCondition == False``) 입력에서 실제로 무엇을 하는지 OBSERVE 한 뒤, 그 동작을
그대로 단언한다. 따라서 미수정 코드에서 PASS(baseline 확립)하고, 수정 후에도 계속 PASS 해야
한다(회귀 가드).

본 파일은 두 종류의 단언으로 구성된다:

  1) **미수정 placement 경로 baseline** (지금 PASS) — ``build_native_cover`` /
     ``build_native_diagram`` / ``_classify_slide_role`` / ``_select_render_plan`` 의
     비버그 입력 좌표·플랜을 관찰·고정한다.

  2) **신규 ``ai_engine.layout_geometry`` no-op 동등성** (모듈 신설 전엔 skip) —
     충돌 회피 기하 함수(``vertical_stack`` / ``resolve_collisions`` /
     ``place_badge_in_gutter`` / ``body_safe_area``)가 겹침 임계 미만 입력에 대해 입력을
     그대로 반환함을 단언한다. ``layout_geometry`` 는 task 3.1 에서 생성되므로, 모듈이 없으면
     ``_HAS_LG`` 가 거짓이 되어 해당 테스트는 자동 skip 되고, 모듈이 생기면(task 3.8 재실행)
     자동 활성화되어 실제 no-op 동등성을 검증한다.

보존 동작(design Preservation Checking / Property 3):

  PRES-1  표지 비겹침 보존 (Req 3.1)
      이미 비겹침인 짧은 제목/부제 박스 → ``vertical_stack`` 결과 == 입력 좌표(바이트 동등).

  PRES-2  배지 거터 보존 (Req 3.1)
      ``place_badge_in_gutter`` 는 배지를 라벨 박스 *밖* 거터에 두어 라벨과의 겹침이 0 이다
      (배지가 이미 라벨 밖이면 라벨 침범 없음 — no-op 동등).

  PRES-3  본문 분리 보존 (Req 3.5)
      백드롭(흰 콘텐츠 패널)으로 이미 분리된 경로의 본문 region 은 그대로 유지된다
      (``build_native_diagram(backdrop=True)`` 의 흰 패널 baseline +
       ``body_safe_area`` 가 백드롭/구워진-텍스트 아님 입력에서 desired 그대로 반환).

  PRES-4  구조형 네이티브 보존 (Req 3.6)
      진짜 구조형(흐름/트리/아키텍처)은 ``_classify_slide_role`` 가 ``structural`` 로 분류하고
      편집 가능 네이티브 도형으로 렌더되며(래스터 이미지 캐리어로 전환되지 않음),
      ``_select_render_plan`` 의 주 렌더러는 ``NATIVE_SHAPES`` 다.

  PRES-5  임계 근방 no-op (Req 3.1, design Property 3)
      ``resolve_collisions`` / ``vertical_stack`` / ``body_safe_area`` 는 겹침 < 임계 입력에
      대해 입력 좌표를 그대로 반환한다.

헤르메틱 — 네트워크 0. 본 테스트는 ``_tool_generate_pptx`` 통합 경로를 구동하지 않고,
순수 결정 함수(``_classify_slide_role`` / ``_select_render_plan`` / ``layout_geometry`` 함수)와
인메모리 네이티브 도형 조립(``build_native_*``)만 구동하므로 게이트웨이/Vertex/HTML 렌더가
호출될 여지가 없다(네트워크 호출 0). 겹침은 기존 audit 도구의 축-정렬 교집합 정의
(``scripts/audit_pptx_textbox_overlap.py`` 의 ``ov``)로 측정한다.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_overlay_collision_preservation_pbt.py -p no:cacheprovider -q

_Preservation: Preservation Requirements 전체 (Req 3.1, 3.5, 3.6) / design Property 3_
_Requirements: 3.1, 3.5, 3.6_
"""
from __future__ import annotations

import os
import sys
import ast
import inspect

# Make ai_engine (repo root) and the audit tools (scripts/) importable.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

# Native assembly path (real code under test) — no network.
from ai_engine.native_diagram_pptx import (  # noqa: E402
    build_native_cover,
    build_native_diagram,
)
# Pure decision functions (real code under test) — no network.
from ai_engine.server import _classify_slide_role, _select_render_plan  # noqa: E402

# EXISTING audit measurement function (reused — same axis-aligned intersection).
import audit_pptx_textbox_overlap as tov  # noqa: E402

# ---------------------------------------------------------------------------
# layout_geometry는 task 3.1에서 생성되는 신규 순수 기하 모듈. 미수정 시점엔 없으므로
# import를 시도해 플래그를 세팅하고, no-op 동등성 테스트는 skipif로 가드한다.
# 모듈이 생기면(task 3.8 재실행) 자동 활성화되어 실제 no-op를 검증한다.
# ---------------------------------------------------------------------------
try:
    import ai_engine.layout_geometry as lg  # noqa: E402
    _HAS_LG = True
except Exception:  # pragma: no cover - 미수정 시점 경로
    lg = None
    _HAS_LG = False

_LG_REASON = "ai_engine.layout_geometry 아직 없음 — task 3.1에서 생성, task 3.8에서 활성화"

# design Bug Condition: 의미있는 겹침 임계(작은 박스 면적의 10%).
THRESHOLD = 0.10
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ===========================================================================
# 공통 헬퍼
# ===========================================================================
def _blank_169_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs):
    # layout[6] = 완전 빈 레이아웃(placeholder 간섭 없음).
    return prs.slides.add_slide(prs.slide_layouts[6])


def _area(box):
    # box = (head, l, t, w, h) — audit 도구와 동일 형태.
    return max(0.0, box[3]) * max(0.0, box[4])


def _pictures(slide):
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                out.append(sh)
        except Exception:
            continue
    return out


def _autoshapes_geom(slide):
    """비-그림 AutoShape 들의 (left,top,width,height) 인치 목록."""
    EMU = 914400.0
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
            if None in (l, t, w, h):
                continue
            out.append((l / EMU, t / EMU, w / EMU, h / EMU))
        except Exception:
            continue
    return out


def _rect_coords(r):
    """layout_geometry Rect를 (left, top, width, height) 튜플로 정규화.

    design은 Rect = (left, top, width, height) 인치 튜플로 정의하지만, 구현이
    namedtuple/객체일 수도 있으므로 양쪽을 모두 수용한다(시그니처 도입에 견고)."""
    if hasattr(r, "left") and hasattr(r, "top") and hasattr(r, "width") and hasattr(r, "height"):
        return (float(r.left), float(r.top), float(r.width), float(r.height))
    seq = list(r)
    return (float(seq[0]), float(seq[1]), float(seq[2]), float(seq[3]))


def _approx_rect(a, b, tol=1e-6):
    ca, cb = _rect_coords(a), _rect_coords(b)
    return all(abs(x - y) <= tol for x, y in zip(ca, cb))


# ===========================================================================
# 입력 풀 (분류가 결정론적이 되도록 큐레이션 — 기존 vertex_images 보존 테스트 컨벤션)
# ===========================================================================
# 진짜 구조형 슬라이드 — _classify_slide_role 가 "structural" 을 반환해야 한다(관찰 확인).
_STRUCTURAL_SLIDES = [
    {"title": "업무 처리 프로세스", "bullets": ["접수 -> 검토 -> 승인 -> 완료"]},
    {"title": "데이터 처리 흐름", "bullets": ["수집 -> 정제 -> 분석 -> 시각화"]},
    {"title": "데이터 처리 프로세스 흐름", "bullets": ["수집 -> 정제 -> 적재 -> 분석 -> 적용"]},
    {"title": "배포 흐름 프로세스", "bullets": ["빌드 -> 테스트 -> 스테이징 -> 배포"]},
]

# 구조형 다이어그램 타입(흐름/트리/아키텍처) — build_native_diagram 네이티브 렌더 대상.
_STRUCTURAL_DIAGRAMS = [
    ("flow", "수집 -> 정제 -> 적재 -> 분석"),
    ("tree", "ROOT\n  자식A\n  자식B\n  자식C"),
    ("architecture", "프론트엔드 계층\n백엔드 계층\n데이터 계층"),
]


# ===========================================================================
# PRES-4 (Req 3.6) — 구조형 네이티브 보존 (지금 PASS — 미수정 baseline)
# ===========================================================================
@settings(max_examples=12, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES))
def test_pres4_structural_renders_native_preserved(slide):
    """진짜 구조형 슬라이드는 structural 로 분류되고 편집 가능 네이티브 도형으로 렌더되며,
    렌더 플랜의 주 렌더러는 NATIVE_SHAPES 다(구워진 이미지 캐리어로 전환되지 않음)."""
    # 관찰·고정 1: role 분류는 structural.
    role = _classify_slide_role(slide, is_cover=False)
    assert role == "structural", f"전제/보존 실패 — 구조형이어야 함 (role={role!r})"

    # 비버그 입력 전제: 배경 이미지 없음(defectB 조건의 bgImage 부재) → 본 슬라이드는
    # 역할/플랜 차원에서 비버그.
    # 관찰·고정 2: 구조형 네이티브 렌더 — 래스터 PICTURE 없음 + 라벨이 편집 가능 텍스트로 보존.
    prs = _blank_169_presentation()
    sl = _blank_slide(prs)
    drew = build_native_diagram(
        sl, "flow", "\n".join(slide["bullets"]),
        region=(0.6, 1.7, 12.1, 5.2), palette=None, title=slide["title"],
    )
    assert drew, "구조형 네이티브 다이어그램이 그려져야 함(보존)"
    assert _pictures(sl) == [], "구조형은 래스터 이미지가 임베드되면 안 됨(편집 가능 네이티브 유지)"

    # 관찰·고정 3: _select_render_plan 의 주 렌더러는 NATIVE_SHAPES, Vertex 손실-0.
    plan = _select_render_plan(
        has_vertex_image=False, has_native_diagram=True,
        has_image_file=False, has_slide_bg=False,
        role="structural", html_enabled=False,
    )
    assert plan["primary"] == "NATIVE_SHAPES", f"구조형 주 렌더러 보존 실패: {plan}"


@settings(max_examples=9, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(spec=st.sampled_from(_STRUCTURAL_DIAGRAMS))
def test_pres4_structural_diagram_no_raster(spec):
    """흐름/트리/아키텍처 구조형은 모두 네이티브 도형으로 렌더되어 래스터 이미지가 없다."""
    dtype, content = spec
    prs = _blank_169_presentation()
    sl = _blank_slide(prs)
    drew = build_native_diagram(sl, dtype, content,
                                region=(0.6, 1.7, 12.1, 5.2), palette=None, title="제목")
    assert drew, f"{dtype} 구조형 네이티브 렌더 실패(보존)"
    assert _pictures(sl) == [], f"{dtype} 구조형에 래스터 이미지가 임베드되면 안 됨(네이티브 유지)"


# ===========================================================================
# PRES-3 (Req 3.5) — 본문 분리 보존: 백드롭(흰 콘텐츠 패널) baseline (지금 PASS)
# ===========================================================================
@settings(max_examples=9, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(spec=st.sampled_from(_STRUCTURAL_DIAGRAMS))
def test_pres3_backdrop_white_panel_separation_preserved(spec):
    """backdrop=True 경로는 본문을 배경과 분리하는 흰 콘텐츠 패널(≈0.5,0.5,12.333,6.58)을
    먼저 깐다 — 본문이 배경 위에 직접 올라가지 않도록 분리하는 baseline. 수정 후에도 이
    분리 동작은 보존되어야 한다(Req 3.5)."""
    dtype, content = spec
    prs = _blank_169_presentation()
    sl = _blank_slide(prs)
    drew = build_native_diagram(sl, dtype, content,
                                region=(0.6, 1.7, 12.1, 5.2), palette=None,
                                title="제목", backdrop=True)
    assert drew, f"{dtype} backdrop 네이티브 렌더 실패(전제)"

    # 관찰·고정: 본문 영역 대부분을 덮는 큰 분리 패널(흰 콘텐츠 카드)이 존재.
    #   build_native_diagram backdrop 분기: cx0=0.5, ctop=0.5, card_w=12.333,
    #   ch=max(2.0,(7.5-0.42)-0.5)=6.58.
    geoms = _autoshapes_geom(sl)
    panel = [g for g in geoms
             if abs(g[0] - 0.5) <= 0.05 and abs(g[1] - 0.5) <= 0.05
             and abs(g[2] - 12.333) <= 0.1 and abs(g[3] - 6.58) <= 0.2]
    assert panel, (
        "백드롭 분리 패널(흰 콘텐츠 카드)이 보존되어야 함 — "
        f"본문이 배경과 분리되지 않음. autoshapes={geoms[:6]}"
    )


# ===========================================================================
# PRES-1 (Req 3.1) — vertical_stack no-op 동등성 (layout_geometry, 모듈 신설 후 활성)
# ===========================================================================
@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
@settings(max_examples=60, deadline=None)
@given(
    first_top=st.floats(min_value=0.0, max_value=2.0),
    left=st.floats(min_value=0.0, max_value=6.0),
    width=st.floats(min_value=0.5, max_value=6.0),
    heights=st.lists(st.floats(min_value=0.3, max_value=1.2), min_size=2, max_size=5),
    seps=st.lists(st.floats(min_value=0.2, max_value=1.0), min_size=1, max_size=5),
    gap=st.floats(min_value=0.0, max_value=0.15),
)
def test_pres1_vertical_stack_noop_when_already_separated(first_top, left, width, heights, seps, gap):
    """이미 (gap 이상으로) 세로 분리된 박스들 → vertical_stack 은 입력 좌표를 그대로 반환."""
    # gap 보다 큰 분리로 비겹침 스택을 구성(첫 박스 top 보존 규칙과도 부합).
    boxes = []
    top = first_top
    for i, h in enumerate(heights):
        boxes.append((left, top, width, h))
        sep = seps[i % len(seps)] + gap + 0.05  # 항상 gap 보다 큰 분리
        top = top + h + sep
    out = lg.vertical_stack([tuple(b) for b in boxes], gap=gap)
    assert len(out) == len(boxes), "vertical_stack 은 입력과 동일 개수를 반환해야 함"
    for inp, res in zip(boxes, out):
        assert _approx_rect(inp, res), (
            f"이미 비겹침인 박스는 좌표 불변이어야 함 — in={inp} out={_rect_coords(res)}"
        )


# ===========================================================================
# PRES-2 (Req 3.1) — place_badge_in_gutter 라벨 밖 배치 (겹침 0) (layout_geometry)
# ===========================================================================
@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
@settings(max_examples=60, deadline=None)
@given(
    lx=st.floats(min_value=1.0, max_value=9.0),
    ly=st.floats(min_value=0.0, max_value=6.0),
    lw=st.floats(min_value=1.0, max_value=5.0),
    lh=st.floats(min_value=0.4, max_value=1.5),
    diameter=st.floats(min_value=0.2, max_value=0.9),
    gap=st.floats(min_value=0.0, max_value=0.2),
)
def test_pres2_place_badge_in_gutter_zero_overlap(lx, ly, lw, lh, diameter, gap):
    """place_badge_in_gutter 는 배지를 라벨 박스 '밖' 좌측 거터에 두어 라벨과의 겹침이 0 이다
    (배지가 라벨을 침범하지 않음 — 거터 보존)."""
    label = (lx, ly, lw, lh)
    badge = lg.place_badge_in_gutter(label, diameter, gutter="left", gap=gap)
    bcoords = _rect_coords(badge)
    # 정사각 diameter 보존.
    assert abs(bcoords[2] - diameter) <= 1e-6, f"배지 폭은 diameter 여야 함: {bcoords}"
    assert abs(bcoords[3] - diameter) <= 1e-6, f"배지 높이는 diameter 여야 함: {bcoords}"
    # 라벨 왼쪽 밖에 위치(우단 <= 라벨 좌단).
    assert bcoords[0] + bcoords[2] <= lx + 1e-6, (
        f"배지는 라벨 왼쪽 거터에 있어야 함 — badge right={bcoords[0]+bcoords[2]:.4f}, label left={lx:.4f}"
    )
    # 라벨과의 겹침 0(audit 정의와 동일한 layout_geometry.overlap_area 사용).
    label5 = ("L", lx, ly, lw, lh)
    badge5 = ("B", bcoords[0], bcoords[1], bcoords[2], bcoords[3])
    assert tov.ov(badge5, label5) <= 1e-9, "배지∩라벨 겹침은 0 이어야 함(거터 보존)"


# ===========================================================================
# PRES-5 (Req 3.1, design Property 3) — 임계 근방 no-op 동등성 (layout_geometry)
# ===========================================================================
@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
@settings(max_examples=60, deadline=None)
@given(
    left=st.floats(min_value=0.0, max_value=4.0),
    width=st.floats(min_value=0.5, max_value=4.0),
    first_top=st.floats(min_value=0.0, max_value=1.0),
    heights=st.lists(st.floats(min_value=0.3, max_value=1.0), min_size=2, max_size=4),
)
def test_pres5_resolve_collisions_noop_when_separated(left, width, first_top, heights):
    """겹침이 임계 미만(완전 분리)인 박스 집합 → resolve_collisions 는 입력 그대로 반환."""
    boxes = []
    top = first_top
    for h in heights:
        boxes.append((left, top, width, h))
        top = top + h + 0.5  # 완전 분리(겹침 0 < 임계)
    out = lg.resolve_collisions([tuple(b) for b in boxes], threshold=THRESHOLD, axis="vertical")
    assert len(out) == len(boxes)
    for inp, res in zip(boxes, out):
        assert _approx_rect(inp, res), (
            f"임계 미만 입력은 좌표 불변이어야 함(no-op) — in={inp} out={_rect_coords(res)}"
        )


@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
@settings(max_examples=40, deadline=None)
@given(
    bx=st.floats(min_value=0.0, max_value=3.0),
    by=st.floats(min_value=1.5, max_value=3.0),
    bw=st.floats(min_value=4.0, max_value=8.0),
    bh=st.floats(min_value=2.0, max_value=4.0),
    has_bg=st.booleans(),
)
def test_pres5_body_safe_area_noop_for_backdrop(bx, by, bw, bh, has_bg):
    """body_safe_area 는 배경이 없거나(백드롭/구워진-텍스트 아님) 분리된 경우 desired 를 그대로
    반환한다(본문 분리 보존, Req 3.5 / design Property 3)."""
    slide = (0.0, 0.0, SLIDE_W_IN, SLIDE_H_IN)
    desired = (bx, by, bw, bh)
    # 백드롭(흰 패널) 또는 배경 없음 — 구워진-텍스트 아님 → desired 보존.
    bg = (0.0, 0.0, SLIDE_W_IN, SLIDE_H_IN) if has_bg else None
    out = lg.body_safe_area(slide, bg, has_baked_text=False, desired=desired)
    assert _approx_rect(desired, out), (
        f"백드롭/구워진-텍스트 아님 입력은 본문 region 보존이어야 함 — "
        f"desired={desired} out={_rect_coords(out)}"
    )


# ==========================================================================
# PROPERTY 4 — Preservation: Vertex 손실-0 불변식 보존 (design Property 4)
#
# _For any_ 슬라이드 미디어 상태에 대해, 본 수정 후에도 ``_select_render_plan`` 은 생성된
# Vertex 이미지를 폐기하지 않는다(``has_vertex_image`` ⇒ ``vertex_slot != "none"``).
# 구조형(role="structural")에서 구워진-텍스트(bg_has_baked_text=True) 이미지를 본문
# 캐리어로 쓰지 않더라도, 생성 이미지는 ``backdrop`` 슬롯으로 보존된다(손실-0). Vertex
# 이미지가 없으면(``has_vertex_image=False``) 슬롯은 의미 없음(``"none"``) — 네이티브/HTML
# 폴백 경로.
#
# 헤르메틱: ``_select_render_plan`` 은 LLM/게이트웨이 호출이 없는 순수 결정 함수다(네트워크
# 0). (has_vertex_image, has_native_diagram, has_image_file, has_slide_bg, role,
# html_enabled, bg_has_baked_text) 조합을 hypothesis 로 광범위 생성해 단언한다.
#
# _Requirements: 3.2, 3.3_
# ==========================================================================
# design classifyRole 반환 열거형 전체(structural 포함).
_ROLES = ["cover", "section", "structural", "content", "visual"]


@settings(max_examples=400, deadline=None)
@given(
    has_vertex_image=st.booleans(),
    has_native_diagram=st.booleans(),
    has_image_file=st.booleans(),
    has_slide_bg=st.booleans(),
    role=st.sampled_from(_ROLES),
    html_enabled=st.booleans(),
    bg_has_baked_text=st.booleans(),
)
def test_prop4_select_render_plan_loss_zero(has_vertex_image, has_native_diagram,
                                            has_image_file, has_slide_bg, role,
                                            html_enabled, bg_has_baked_text):
    """모든 슬라이드 미디어 상태에서 _select_render_plan 의 손실-0 불변식이 보존된다.

    **Validates: Requirements 3.2, 3.3**
    """
    plan = _select_render_plan(
        has_vertex_image=has_vertex_image,
        has_native_diagram=has_native_diagram,
        has_image_file=has_image_file,
        has_slide_bg=has_slide_bg,
        role=role,
        html_enabled=html_enabled,
        bg_has_baked_text=bg_has_baked_text,
    )
    # 반환 계약: 세 키가 모두 존재하고 primary 는 알려진 렌더러.
    assert {"primary", "vertex_slot", "body_separated"} <= set(plan.keys()), plan
    assert plan["primary"] in ("HTML", "NATIVE_SHAPES", "VERTEX_IMAGE"), plan
    assert isinstance(plan["body_separated"], bool), plan

    if has_vertex_image:
        # 손실-0 핵심 불변식: 생성된 Vertex 이미지는 어떤 분기에서도 폐기되지 않는다.
        assert plan["vertex_slot"] != "none", (
            f"손실-0 위반 — has_vertex_image=True인데 vertex_slot='none': {plan}"
        )
        assert plan["vertex_slot"] in ("hero", "backdrop", "visual"), plan
        # 구조형 + 구워진-텍스트: 생성 이미지를 본문 캐리어(visual/hero)로 쓰지 않고
        # NATIVE_SHAPES 를 주 렌더러로 유지하며 이미지는 backdrop 으로 보존(손실-0).
        if role == "structural" and bg_has_baked_text:
            assert plan["primary"] == "NATIVE_SHAPES", (
                f"구조형+구워진텍스트는 NATIVE_SHAPES 가 주 렌더러여야 함: {plan}"
            )
            assert plan["vertex_slot"] == "backdrop", (
                f"구조형+구워진텍스트는 생성 이미지를 backdrop 으로 보존해야 함: {plan}"
            )
            assert plan["body_separated"] is True, (
                f"구조형+구워진텍스트는 본문/배경 분리 신호를 세워야 함: {plan}"
            )
    else:
        # Vertex 이미지가 없으면 슬롯은 의미 없음("none") — 네이티브/HTML 폴백 경로.
        assert plan["vertex_slot"] == "none", (
            f"has_vertex_image=False 면 vertex_slot 은 'none' 이어야 함: {plan}"
        )


# ==========================================================================
# PROPERTY 5 — Preservation: 게이트웨이 제약 보존 (design Property 5)
#
# 신규 충돌 회피 기하 모듈(ai_engine.layout_geometry)은 순수 계산이며 어떤 네트워크/모델/
# 게이트웨이 호출도 하지 않는다. design Property 5: "신규 충돌 회피 기하 함수는 순수 계산
# 이며 어떤 네트워크/모델 호출도 하지 않는다." 이를 import 그래프 + 소스 정적 검사(ast)로
# 단언한다.
#
#   PROP5-A  layout_geometry 가 표준 라이브러리(+__future__/typing)만 import 함
#            (requests/httpx/boto3/ai_engine.* 등 비표준/1st-party import 0).
#   PROP5-B  layout_geometry 소스(식별자/속성/import)에 네트워크·게이트웨이·모델 호출
#            토큰(requests/httpx/boto3/get_vertex/_get_gw/_call_bridge 등)이 없음.
#
# 헤르메틱 — 정적 검사뿐, 네트워크 호출 0.
#
# _Requirements: 3.4_
# ==========================================================================
# 게이트웨이/네트워크/모델 호출을 시사하는 금지 토큰(코드 식별자로 등장하면 안 됨).
_FORBIDDEN_NETWORK_TOKENS = [
    "requests", "httpx", "urllib", "urllib3", "socket", "boto3", "botocore",
    "aiohttp", "vertex", "get_vertex", "get_vertex_image_client", "_get_gw",
    "_call_bridge", "bedrock", "openai", "anthropic", "subprocess",
    "_specialized_model_for_task", "generate",
]


def _layout_geometry_import_tops():
    """layout_geometry 가 import 하는 최상위 모듈 이름 목록을 ast 로 수집."""
    src = inspect.getsource(lg)
    tree = ast.parse(src)
    tops = []
    for node in ast.walk(tree):
        if isinstance(node, ast.Import):
            for alias in node.names:
                tops.append(alias.name.split(".")[0])
        elif isinstance(node, ast.ImportFrom):
            # 상대 import(level>0)는 1st-party 의존 — 순수 모듈엔 없어야 함.
            assert node.level == 0, (
                f"layout_geometry 에 상대 import(1st-party 의존): {ast.dump(node)}"
            )
            if node.module:
                tops.append(node.module.split(".")[0])
    return tops


@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
def test_prop5a_layout_geometry_imports_only_stdlib():
    """PROP5-A: layout_geometry 는 표준 라이브러리(+__future__/typing)만 import 한다.

    **Validates: Requirements 3.4**
    """
    tops = _layout_geometry_import_tops()
    assert tops, "import 문이 하나도 안 잡힘 — 소스 파싱 오류 의심"
    stdlib = getattr(sys, "stdlib_module_names", None)
    for mod in tops:
        # 1st-party(ai_engine) 의존 금지 — 네트워크/게이트웨이 전이 의존 차단.
        assert mod != "ai_engine", (
            f"layout_geometry 가 ai_engine 을 import 함({mod!r}) — 순수 모듈이어야 함"
        )
        if stdlib is not None:
            assert mod == "__future__" or mod in stdlib, (
                f"layout_geometry 가 비표준 모듈을 import 함: {mod!r} — "
                "순수 모듈은 표준 라이브러리만 import 해야 함"
            )


@pytest.mark.skipif(not _HAS_LG, reason=_LG_REASON)
def test_prop5b_layout_geometry_source_has_no_network_calls():
    """PROP5-B: layout_geometry 코드에 네트워크/게이트웨이/모델 호출 토큰이 없다.

    주석/docstring 의 설명 텍스트(예: 'gateway', 'vertex' 단어)와 충돌하지 않도록, 문자열
    매칭이 아니라 ast 로 코드 식별자(Name)/속성 접근(Attribute)/import 이름만 수집해 검사
    한다 — 실제 호출 가능한 심볼만 대상으로 삼는다.

    **Validates: Requirements 3.4**
    """
    src = inspect.getsource(lg)
    tree = ast.parse(src)
    code_names = set()
    for node in ast.walk(tree):
        if isinstance(node, ast.Name):
            code_names.add(node.id)
        elif isinstance(node, ast.Attribute):
            code_names.add(node.attr)
        elif isinstance(node, ast.Import):
            for alias in node.names:
                code_names.add(alias.name.split(".")[0])
        elif isinstance(node, ast.ImportFrom):
            if node.module:
                code_names.add(node.module.split(".")[0])
            for alias in node.names:
                code_names.add(alias.name.split(".")[0])
    hits = sorted(t for t in _FORBIDDEN_NETWORK_TOKENS if t in code_names)
    assert not hits, (
        f"layout_geometry 코드에 네트워크/게이트웨이/모델 호출 토큰이 있음: {hits} — "
        "기하 함수는 순수 계산이어야 한다(design Property 5)"
    )


# ==========================================================================
# PROPERTY 18 — 겹침 0 산출물 (spec: pptx-ultra-quality-hybrid-render, task 8.5)
#
# design Property 18 (Validates: Requirements 4.6, 5.2):
#   "*For any* 하이브리드 렌더로 생성한 슬라이드에 대해, 도형-도형 및 텍스트박스-텍스트박스
#    겹침 면적은 0 EMU이며 audit_pptx_textbox_overlap.main(path)가 판정한 겹침 면적
#    0.05in² 초과 텍스트박스 쌍 개수는 0이다."
#
# 본 테스트는 실제 `_tool_generate_pptx`를 `AE_HYBRID_RENDER=1` + Vertex 비활성 + HTML off로
# 헤르메틱(네트워크 0) 구동해 실 .pptx를 생성한 뒤, (a) 모든 슬라이드의 겹침-검사 대상
# (텍스트 보유) 도형쌍 겹침 면적이 0 EMU이고, (b) `audit_pptx_textbox_overlap.main(path)`가
# 보고하는 0.05in² 초과 텍스트박스 쌍 개수가 0임을 100+ iteration 으로 검증한다.
#
# 겹침 스코프(정확성): "도형-도형 및 텍스트박스-텍스트박스"는 겹침-검사 대상 도형
#   = 텍스트 보유 비-PICTURE 도형(TEXT_BOX + 텍스트 보유 AUTO_SHAPE) 으로 측정한다. 이는
#   기존 감사 도구 `audit_pptx_textbox_overlap.boxes` 의 추출 정의 및
#   `native_layout_renderer._participates_in_collision`(텍스트 없는 장식 배경은 의도적
#   레이어링으로 제외) 과 정확히 일치한다. 즉 텍스트 없는 장식 배경 컨테이너는 제외하고,
#   실제 콘텐츠를 나르는 도형쌍의 겹침만 0 EMU 인지 본다(감사 도구와 동일 기준).
#
# 입력 풀: role="content" 로 결정론 분류되는 **순수 텍스트** content 슬라이드만 사용한다.
#   이 입력군은 비하이브리드(카드/네이티브 폴백) 기준선에서 겹침 0(main 초과쌍 0)이 실측으로
#   성립하므로, Property 18 은 올바른 구현이 달성 가능한 정당한 목표다(unsatisfiable 아님).
#
# 헤르메틱: Vertex disabled 스텁(get_vertex_image_client), _call_bridge/_tool_generate_image
#   목, Chrome 경로(_find_local_chrome) 목 → 네트워크 호출 0.
#
# _Requirements: 4.6, 5.2_
# ==========================================================================
import io as _io18  # noqa: E402
import json as _json18  # noqa: E402
import asyncio as _asyncio18  # noqa: E402
import tempfile as _tempfile18  # noqa: E402
import contextlib as _contextlib18  # noqa: E402
from unittest.mock import patch as _patch18  # noqa: E402

import ai_engine.server as _srv18  # noqa: E402
import ai_engine.vertex_image_module as _vim18  # noqa: E402


class _P18DisabledVertexClient:
    """Vertex 비활성 스텁 — enabled=False. generate 는 네트워크 없이 즉시 실패 반환하며
    호출 횟수를 기록해(정상 경로에서는 0회) 네트워크 0을 단언할 수 있게 한다."""

    enabled = False

    def __init__(self) -> None:
        self.calls = 0

    async def generate(self, *a, **k):  # pragma: no cover - 정상 경로에선 미호출
        self.calls += 1
        return {"error": "test-disabled"}


async def _p18_img_gen_disabled(*_a, **_k):
    """_tool_generate_image 목 — 경로 미반환(네트워크 0)."""
    return _json18.dumps({"error": "test-disabled"})


# role="content" 로 결정론 분류되는 순수 텍스트 content 슬라이드(비하이브리드 기준선 겹침 0).
_P18_CONTENT_SLIDES = [
    {"title": "환영합니다", "bullets": ["반갑습니다", "함께 잘 부탁드립니다", "좋은 하루 되세요"]},
    {"title": "팀 인사", "bullets": ["좋은 하루입니다", "즐겁게 시작해요", "화이팅"]},
    {"title": "오늘의 메시지", "bullets": ["감사합니다", "끝까지 함께 가요"]},
    {"title": "마무리 인사", "bullets": ["수고하셨습니다", "다음에 또 만나요", "건강하세요"]},
]


def _run_hybrid_render_html_off(slides):
    """AE_HYBRID_RENDER=1 + Vertex 비활성 + HTML off 로 실 .pptx 를 헤르메틱 생성한다.

    네트워크 0: Vertex disabled 스텁 + _call_bridge/_tool_generate_image/Chrome 목.
    반환: 생성된 .pptx 절대경로."""
    fake = _P18DisabledVertexClient()
    proj = _tempfile18.mkdtemp()
    env = {
        "AE_HYBRID_RENDER": "1",             # 하이브리드 라우팅 활성(본 프로퍼티 대상)
        "AE_ENABLE_HTML_SLIDES": "0",        # HTML off → Chrome 불필요(네이티브/편집 경로)
        "AE_ENABLE_VERTEX_IMAGE": "0",       # Vertex 비활성(헤르메틱)
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # 결정론: 게이트웨이 구조화 없음
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    with _patch18.dict(os.environ, env, clear=False), \
            _patch18.object(_vim18, "get_vertex_image_client", lambda **_k: fake), \
            _patch18.object(_srv18, "_call_bridge", lambda *a, **k: None), \
            _patch18.object(_srv18, "_find_local_chrome", lambda: ""), \
            _patch18.object(_srv18, "_tool_generate_image", _p18_img_gen_disabled):
        raw = _asyncio18.run(_srv18._tool_generate_pptx(
            {"title": "겹침0 검증 덱", "slides": [dict(s) for s in slides]},
            project_path=proj))
    result = _json18.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    # Vertex 비활성 → 이미지 생성 경로 진입 없음(네트워크 0 재확인).
    assert fake.calls == 0, f"Vertex 비활성인데 generate 호출됨(헤르메틱 위반): {fake.calls}"
    return result["absPath"]


def _max_participating_overlap(slide):
    """겹침-검사 대상(텍스트 보유) 도형쌍의 최대 겹침 면적(in²)과 쌍 라벨.

    기존 감사 도구 ``audit_pptx_textbox_overlap.boxes``(텍스트 보유 비-PICTURE 도형 추출) +
    ``ov``(축-정렬 교집합) 를 그대로 재사용한다 — 텍스트 없는 장식 배경은 제외(감사/
    _participates_in_collision 와 동일 기준)."""
    bs = tov.boxes(slide)
    worst = 0.0
    pair = None
    for i in range(len(bs)):
        for j in range(i + 1, len(bs)):
            a = tov.ov(bs[i], bs[j])
            if a > worst:
                worst = a
                pair = (bs[i][0], bs[j][0])
    return worst, pair


def _main_over_threshold_pairs(path):
    """``audit_pptx_textbox_overlap.main(path)`` 의 stdout 를 캡처해 0.05in² **초과** 텍스트박스
    쌍 개수를 센다.

    main 은 초과 쌍마다 ``... 겹침 X.XXin² ...`` 라인을, 초과 쌍이 없는 슬라이드엔
    ``텍스트 박스 간 겹침 없음`` 라인을 출력한다. 초과 쌍 라인만 ``in²`` 토큰을 포함하므로
    그 라인 수가 곧 초과 쌍 개수다(요약 라인 없음)."""
    buf = _io18.StringIO()
    with _contextlib18.redirect_stdout(buf):
        tov.main(path)
    out = buf.getvalue()
    return sum(1 for ln in out.splitlines() if "in²" in ln)


@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_P18_CONTENT_SLIDES))
def test_prop18_hybrid_render_zero_overlap(slide):
    """Property 18: 하이브리드 렌더 산출물의 도형-도형/텍스트박스-텍스트박스 겹침 0.

    하이브리드 렌더(AE_HYBRID_RENDER=1, Vertex 비활성, 헤르메틱)로 생성한 .pptx 의 모든
    슬라이드에서 (a) 겹침-검사 대상(텍스트 보유) 도형쌍 겹침 면적이 0 EMU 이고, (b)
    ``audit_pptx_textbox_overlap.main(path)`` 가 보고하는 0.05in² 초과 텍스트박스 쌍 개수가
    0 이다.

    **Validates: Requirements 4.6, 5.2**
    """
    path = _run_hybrid_render_html_off([slide])

    # (b) 감사 도구 판정 — 0.05in² 초과 텍스트박스 쌍 0개 (R5.2).
    over = _main_over_threshold_pairs(path)
    assert over == 0, (
        f"audit_pptx_textbox_overlap.main 이 0.05in² 초과 텍스트박스 쌍 {over}개 보고 — "
        f"겹침0 산출물 위반 (Property 18 / R5.2)")

    # (a) 슬라이드별 겹침-검사 대상 도형쌍 겹침 면적 0 EMU (R4.6).
    prs = Presentation(path)
    for idx, sl in enumerate(prs.slides, 1):
        worst, pair = _max_participating_overlap(sl)
        assert worst <= 1e-9, (
            f"[슬라이드 {idx}] 도형-도형/텍스트박스-텍스트박스 겹침 {worst:.4f}in² (>0 EMU) — "
            f"{pair} (Property 18 / R4.6)")


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
