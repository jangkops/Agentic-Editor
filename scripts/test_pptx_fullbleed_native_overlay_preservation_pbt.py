"""Preservation property tests — spec: pptx-fullbleed-native-overlay-collision
(bugfix), Task 2.

PROPERTY 2 — Preservation (비결함 입력 동작 보존).

이 PBT는 풀블리드-네이티브 오버레이 충돌/경계밖 수정(tasks 3.x)이 **변경해서는 안 되는**
동작을 고정한다. **observation-first** 방법론을 따른다 — 각 단언은 먼저 *미수정* 코드
(``ai_engine/server.py`` 의 ``_tool_generate_pptx`` 결정 seam)가 비결함
(``isBugCondition == False``) 입력에서 실제로 무엇을 하는지 OBSERVE 한 뒤 그 동작을 그대로
단언한다. 따라서 미수정 코드에서 PASS(baseline 확립)하고 수정 후에도 계속 PASS 해야 한다
(회귀 가드).

design "Preservation Requirements" / "Preservation Checking" / Property 2 를 따라, 결함
조건이 아닌 다음 입력군을 다룬다:

  PRES-1  직접 네이티브 경로 보존 (Req 3.1)
      slideBackground 미설정 + 네이티브 다이어그램(흐름/트리/아키텍처) 슬라이드는 콘텐츠가
      구워진 풀블리드 PICTURE 없이 편집 가능 네이티브 도형으로 렌더된다. 본문 콘텐츠
      슬라이드에 (a) 풀블리드 구워진 배경이 없고(따라서 텍스트↔구워진배경 겹침이 구조적으로
      0), (b) 경계(0,0,13.333,7.5) 밖 도형이 없으며, (c) 라벨이 편집 가능 텍스트로 남는다.
      이 경로는 design 이 "현재 정상(clean)" 이라고 명시한 경로 — 그대로 잠근다.

  PRES-2  손실-0 임베드 보존 (Req 3.2)
      Vertex 이미지가 생성되면(비구조 visual/content) 생성된 모든 이미지가 ``ppt/media/*``
      에 임베드된다(unused == 0) — pptx-quality-vertex-images 손실-0 불변식.

  PRES-3  명시 imageFile / 장식 slideBackground 우선순위 보존 (Req 3.3)
      caller 가 명시한 imageFile/slideBackground 가 주 렌더러로 유지된다: Vertex 사전생성은
      그 슬라이드를 건너뛰고 caller 의 정확한 바이트가 임베드된다. 결정 함수
      ``_select_render_plan`` 의 분기(imageFile→VERTEX_IMAGE, slideBackground→HTML)도 보존.

  PRES-4  풀블리드 없는 네이티브 텍스트 슬라이드 레이아웃/여백 보존 (Req 3.4)
      풀블리드 배경 없는 비구조 텍스트 슬라이드는 기존 레이아웃·여백을 유지한다: 풀블리드
      PICTURE 없음, 본문 불릿 텍스트 보존, 텍스트 도형이 경계 안 + 좌측 여백 유지.

  PRES-5  비결함 랜덤 덱 보존 (Req 3.5)
      구조형/콘텐츠/비주얼이 무작위로 섞인 비결함 덱에서 슬라이드 수(표지+N), 손실-0(unused
      ==0), 본문 슬라이드 경계 안이 보존된다(D1 풀블리드 1회·z-order 불변식 영역).

모든 테스트는 hermetic — 네트워크 0. Bedrock 게이트웨이, Vertex 클라이언트
(``get_vertex_image_client`` / ``generate``), HTML→PNG 렌더러, ``_tool_generate_image`` 를
모두 mock 한다. 모든 파이프라인 드라이버는 HTML OFF(Chrome 불필요)로 구동하며
``_find_local_chrome`` 를 "" 로 패치해 Chrome 의존을 제거한다(hang 방지). 생성된 ``.pptx`` 를
zip 으로 열어 ``ppt/media/*`` 바이트 + 도형 기하를 검사한다. 겹침/경계밖 측정은 기존 감사
스크립트(``scripts/audit_pptx_textbox_overlap.py`` 의 ``ov``,
``scripts/audit_pptx_zorder_break.py`` 의 풀블리드/경계 정의)와 동일한 축-정렬 교집합 정의를
재사용한다.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_fullbleed_native_overlay_preservation_pbt.py -p no:cacheprovider -q

_Requirements: 3.1, 3.2, 3.3, 3.4, 3.5_
"""
from __future__ import annotations

import io
import os
import sys
import json
import base64
import asyncio
import zipfile
import tempfile
from unittest.mock import patch

# Make ai_engine (repo root) and the audit tools (scripts/) importable.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402

# EXISTING audit measurement function (reused — same axis-aligned intersection).
import audit_pptx_textbox_overlap as tov  # noqa: E402

EMU = 914400.0
SW, SH = 13.333, 7.5
OVERLAP_THRESHOLD = 0.10   # design Bug Condition — 의미있는 겹침 임계(텍스트 면적 10%)
BOUNDS_EPS = 0.05          # 경계 검사 허용 오차(audit_pptx_zorder_break 와 정합)


# --------------------------------------------------------------------------
# Formal spec mirror — isBugCondition (design "Formal Specification").
# Used as a PRECONDITION guard at the SCENARIO level: every input these tests
# drive is constructed to be a non-bug input (no baked-content full-bleed
# overlapping native text, no off-slide shapes), so the preservation assertions
# are about the unchanged (non-bug) domain.
# --------------------------------------------------------------------------
def is_bug_condition_scenario(state: dict) -> bool:
    """design isBugCondition 의 시나리오 수준 미러.

    collision := 콘텐츠가 구워진 풀블리드 배경 + 같은 콘텐츠 네이티브 레이어 겹침 ≥ 임계.
    offSlide  := 슬라이드 경계 밖 도형 존재.
    (구워진-콘텐츠 여부는 픽셀이 아니라 시나리오 구성으로 결정 — caller 장식 배경은
     baked 가 아니므로 collision 에 기여하지 않는다.)
    """
    collision = state.get("hasBakedFullbleed", False) and state.get("hasNativeOverlap", False)
    off_slide = state.get("hasOffSlide", False)
    return collision or off_slide


# --------------------------------------------------------------------------
# Hermetic fakes (mirror the existing pptx-quality-vertex-images test files)
# --------------------------------------------------------------------------
def _make_png(tag: int) -> bytes:
    """Produce a unique, valid PNG so each generated image is byte-distinct."""
    img = Image.new(
        "RGB",
        (40 + (tag % 11), 30 + (tag % 7)),
        (tag % 256, (tag * 7) % 256, (tag * 31) % 256),
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class _FakeVertexClient:
    """Stand-in for VertexImageClient — 'enabled', returns unique PNGs and records
    the raw bytes so the test can verify loss-zero (every generated image embedded)."""

    def __init__(self, enabled: bool = True) -> None:
        self.enabled = enabled
        self.calls = 0
        self.generated_raw: list[bytes] = []

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.calls += 1
        raw = _make_png(2000 + self.calls)
        self.generated_raw.append(raw)
        return {"images": [base64.b64encode(raw).decode("ascii")]}


async def _img_gen_disabled(*_a, **_k):
    """Stand-in for _tool_generate_image — never returns a path (no network)."""
    return json.dumps({"error": "test-disabled"})


# --------------------------------------------------------------------------
# HTML-OFF pipeline driver (no Chrome needed) — exercises the real decision seam.
# --------------------------------------------------------------------------
def _run_html_off(slides, *, vertex_enabled: bool, extra_env=None):
    """Drive the real _tool_generate_pptx with HTML OFF. Returns (fake, result)."""
    fake = _FakeVertexClient(enabled=vertex_enabled)
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",        # HTML off → no Chrome, native/Vertex path
        "AE_PREFER_VERTEX_IMAGE": "1" if vertex_enabled else "0",
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # deterministic: no gateway structuring
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",      # keep cover simple (cover bug is task 1's scope)
        "AE_GENERATED_ROOT": proj,
    }
    if extra_env:
        env.update(extra_env)
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(
            {"title": "보존 검증 덱", "slides": list(slides)}, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return fake, result


# --------------------------------------------------------------------------
# PPTX inspection helpers (reuse audit definitions)
# --------------------------------------------------------------------------
def _pptx_media_bytes(pptx_abs: str) -> list[bytes]:
    with zipfile.ZipFile(pptx_abs) as z:
        return [z.read(n) for n in z.namelist() if n.startswith("ppt/media/")]


def _vertex_usage(fake: _FakeVertexClient, pptx_abs: str):
    """Return (generated, embedded, unused) counts for the Vertex images."""
    media = _pptx_media_bytes(pptx_abs)
    embedded = sum(1 for raw in fake.generated_raw if raw in media)
    generated = len(fake.generated_raw)
    return generated, embedded, generated - embedded


def _slides(pptx_abs: str):
    return list(Presentation(pptx_abs).slides)


def _rect_in(sh):
    """Shape (left, top, width, height) in inches, or None."""
    try:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        if None in (l, t, w, h):
            return None
        return (l / EMU, t / EMU, w / EMU, h / EMU)
    except Exception:
        return None


def _is_fullbleed(r) -> bool:
    """audit_pptx_zorder_break._fullbleed 정의와 동일."""
    return bool(r) and r[0] <= 0.3 and r[1] <= 0.3 and r[2] >= SW * 0.92 and r[3] >= SH * 0.92


def _pictures(slide):
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                out.append(sh)
        except Exception:
            continue
    return out


def _fullbleed_pictures(slide):
    return [sh for sh in _pictures(slide) if _is_fullbleed(_rect_in(sh))]


def _text_shapes(slide):
    """(text, rect) for non-picture shapes carrying text."""
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            pass
        try:
            if getattr(sh, "has_text_frame", False) and (sh.text_frame.text or "").strip():
                r = _rect_in(sh)
                if r is not None:
                    out.append(((sh.text_frame.text or "").strip(), r))
        except Exception:
            continue
    return out


def _all_shape_text(slide) -> str:
    return " ".join(t for t, _ in _text_shapes(slide))


def _within_bounds(r) -> bool:
    """rect 가 슬라이드 경계(0,0,SW,SH) 안인지(audit zorder 의 off-slide 정의와 정합)."""
    if not r:
        return True
    return (r[0] >= -BOUNDS_EPS and r[1] >= -BOUNDS_EPS
            and r[0] + r[2] <= SW + BOUNDS_EPS and r[1] + r[3] <= SH + BOUNDS_EPS)


def _off_slide_shapes(slide):
    """경계 밖 도형의 rect 목록."""
    out = []
    for sh in slide.shapes:
        r = _rect_in(sh)
        if r is not None and not _within_bounds(r):
            out.append(r)
    return out


def _max_text_fullbleed_overlap_ratio(slide) -> float:
    """텍스트 도형 ↔ 풀블리드 PICTURE 최대 면적 겹침 비율(텍스트 면적 기준).

    audit_pptx_textbox_overlap.ov 와 동일한 축-정렬 교집합을 사용한다."""
    fbs = [_rect_in(p) for p in _fullbleed_pictures(slide)]
    worst = 0.0
    for _txt, tr in _text_shapes(slide):
        ta = max(0.0, tr[2]) * max(0.0, tr[3])
        if ta <= 0:
            continue
        t5 = ("T", tr[0], tr[1], tr[2], tr[3])
        for fr in fbs:
            p5 = ("P", fr[0], fr[1], fr[2], fr[3])
            ratio = tov.ov(t5, p5) / ta
            if ratio > worst:
                worst = ratio
    return worst


# --------------------------------------------------------------------------
# Input pools (curated so classification is deterministic — mirror existing
# pptx-quality-vertex-images preservation tests).
# --------------------------------------------------------------------------
# Structural slides — _classify_section_diagram MUST return a real structural
# kind (flow / tree / architecture).
_STRUCTURAL_SLIDES = [
    ("업무 처리 프로세스", ["접수", "검토", "승인", "완료"]),
    ("데이터 처리 흐름", ["수집", "정제", "분석", "시각화"]),
    ("처리 단계", ["요청 접수", "데이터 검증", "결과 반환"]),
    ("시스템 아키텍처", ["프론트엔드 계층", "백엔드 계층", "데이터 계층"]),
    ("배포 흐름 프로세스", ["빌드", "테스트", "스테이징", "배포"]),
]

# Non-structural content slides — _classify_section_diagram MUST return "" so the
# slide takes the plain content / native-text path (NOT a native diagram).
_CONTENT_SLIDES = [
    ("환영합니다", ["반갑습니다", "함께 잘 부탁드립니다"]),
    ("팀 인사", ["좋은 하루입니다", "즐겁게 시작해요"]),
    ("오늘의 메시지", ["감사합니다", "끝까지 함께 가요"]),
    ("마무리 인사", ["수고하셨습니다", "다음에 또 만나요"]),
]

# Visual-intent slides — non-structural title/bullets + a photographic/illustrative
# imagePrompt so _classify_slide_role => "visual" and the Vertex image path runs.
_VISUAL_SLIDES = [
    {"title": "회사 소개", "bullets": ["신뢰", "혁신"],
     "imagePrompt": "a high quality professional photograph of a modern corporate office, natural light"},
    {"title": "브랜드 비전", "bullets": ["미래 지향"],
     "imagePrompt": "a cinematic photograph of a sunrise over a city skyline, warm tones"},
    {"title": "고객 경험", "bullets": ["만족"],
     "imagePrompt": "an editorial illustration of happy customers using a product, flat design"},
]


# ==========================================================================
# PRES-1 (Req 3.1) — 직접 네이티브 경로(slideBackground 미설정)는 깨끗(clean)
# ==========================================================================
@settings(max_examples=10, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES))
def test_pres1_direct_native_path_clean_preserved(slide):
    """직접 네이티브 경로(구조형, slideBackground 미설정)는 구워진 풀블리드 배경이 없고
    경계 밖 도형이 없으며 라벨이 편집 가능 텍스트로 남는다(design 3.1 — 현재 정상 경로)."""
    title, bullets = slide

    # 전제: 진짜 구조형으로 분류 → 직접 네이티브 다이어그램 경로(풀블리드 배경 미생성).
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind in ("flow", "tree", "architecture"), (
        f"전제 실패: {title!r}는 구조형으로 분류돼야 함 (kind={kind!r})")
    # 전제: 비결함 입력(구워진 풀블리드 없음 → collision 불가, 경계밖 없음).
    assert not is_bug_condition_scenario({
        "hasBakedFullbleed": False, "hasNativeOverlap": False, "hasOffSlide": False})

    fake, result = _run_html_off([{"title": title, "bullets": list(bullets)}],
                                 vertex_enabled=False)
    slides = _slides(result["absPath"])
    assert len(slides) >= 2, "표지 + 콘텐츠 슬라이드가 있어야 함"
    content = slides[1]

    # (a) 콘텐츠가 구워진 풀블리드 PICTURE 없음 → 텍스트↔구워진배경 충돌 구조적 0.
    assert _fullbleed_pictures(content) == [], (
        "직접 네이티브 경로 콘텐츠 슬라이드에 풀블리드 배경이 깔리면 안 됨(편집 가능 네이티브 유지)")
    assert _max_text_fullbleed_overlap_ratio(content) < OVERLAP_THRESHOLD, (
        "직접 네이티브 경로: 텍스트↔풀블리드 겹침은 임계 미만(0)이어야 함")
    # 래스터 이미지 자체가 없어야 함(구조형은 편집 가능 도형).
    assert _pptx_media_bytes(result["absPath"]) == [], (
        "구조형 직접 네이티브 경로에는 래스터 이미지가 임베드되면 안 됨")
    # (b) 경계 밖 도형 없음(콘텐츠 슬라이드).
    assert _off_slide_shapes(content) == [], (
        f"직접 네이티브 경로 콘텐츠 슬라이드에 경계 밖 도형이 없어야 함 — {_off_slide_shapes(content)}")
    # (c) 라벨이 편집 가능 텍스트로 보존.
    texts = _all_shape_text(content)
    hits = [b for b in bullets if b in texts]
    assert len(hits) >= max(2, len(bullets) - 1), (
        f"구조형 라벨이 편집 가능 도형 텍스트로 남아야 함 — hits={hits}, bullets={bullets}")


# ==========================================================================
# PRES-2 (Req 3.2) — Vertex 손실-0 임베드 보존
# ==========================================================================
@settings(max_examples=9, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_VISUAL_SLIDES))
def test_pres2_vertex_loss_zero_preserved(slide):
    """비주얼 슬라이드에서 생성된 Vertex 이미지는 모두 ppt/media 에 임베드된다(unused==0)."""
    # 전제: 비주얼(이미지) 작업 → Vertex 사전생성 경로.
    role = srv._classify_slide_role(dict(slide), False, "doc")
    assert role == "visual", f"전제 실패: role={role!r} (visual 기대)"

    fake, result = _run_html_off([dict(slide)], vertex_enabled=True)
    generated, embedded, unused = _vertex_usage(fake, result["absPath"])
    assert generated >= 1, (
        f"비주얼 슬라이드는 Vertex 이미지를 생성해야 함 (generated={generated})")
    assert unused == 0, (
        f"손실-0 위반: 생성된 Vertex 이미지가 폐기됨 "
        f"(generated={generated}, embedded={embedded}, unused={unused})")


# ==========================================================================
# PRES-3 (Req 3.3) — 명시 imageFile / 장식 slideBackground 우선순위 보존
# ==========================================================================
def _write_png(path: str, w: int, h: int, color) -> bytes:
    Image.new("RGB", (w, h), color).save(path, format="PNG")
    with open(path, "rb") as f:
        return f.read()


@settings(max_examples=8, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       which=st.sampled_from(["imageFile", "slideBackground"]))
def test_pres3_caller_specified_image_precedence_preserved(slide, which):
    """caller 가 명시한 imageFile/slideBackground 는 주 렌더러로 유지된다: Vertex 는 그
    슬라이드를 건너뛰고 caller 의 정확한 바이트가 임베드된다(우선순위 보존)."""
    title, bullets = slide
    proj = tempfile.mkdtemp()
    png_path = os.path.join(proj, "caller.png")
    if which == "imageFile":
        caller_bytes = _write_png(png_path, 1200, 700, (11, 22, 33))
    else:
        caller_bytes = _write_png(png_path, 1920, 1080, (7, 8, 9))
    sd = {"title": title, "bullets": list(bullets), which: png_path}

    # 전제: caller 장식/사진 이미지(구워진 콘텐츠 아님) → 비결함.
    assert not is_bug_condition_scenario({
        "hasBakedFullbleed": False, "hasNativeOverlap": False, "hasOffSlide": False})

    fake = _FakeVertexClient(enabled=True)
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex eligible — but must skip this slide
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": [sd]},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"

    # caller 가 이미 imageFile/slideBackground 를 준 슬라이드는 Vertex 사전생성을 건너뛴다.
    assert fake.calls == 0, (
        f"caller가 {which}를 지정한 슬라이드는 Vertex 생성을 건너뛰어야 함 (calls={fake.calls})")
    # caller 의 정확한 바이트가 임베드된다(우선순위 보존).
    media = _pptx_media_bytes(result["absPath"])
    assert caller_bytes in media, (
        f"caller가 지정한 {which} 이미지가 그대로 임베드돼야 함(우선순위 보존)")


def test_pres3_select_render_plan_precedence_preserved():
    """결정 함수 _select_render_plan 의 명시-이미지 우선순위 분기 보존(순수 함수, 네트워크 0).

    - slideBackground 존재  → primary=HTML       (장식 풀블리드가 주 렌더러)
    - imageFile 존재        → primary=VERTEX_IMAGE(caller 명시 이미지가 슬라이드 비주얼)
    """
    bg_plan = srv._select_render_plan(
        has_vertex_image=False, has_native_diagram=False,
        has_image_file=False, has_slide_bg=True,
        role="content", html_enabled=False)
    assert bg_plan["primary"] == "HTML", f"장식 slideBackground 주 렌더러 보존 실패: {bg_plan}"

    img_plan = srv._select_render_plan(
        has_vertex_image=False, has_native_diagram=False,
        has_image_file=True, has_slide_bg=False,
        role="content", html_enabled=False)
    assert img_plan["primary"] == "VERTEX_IMAGE", f"명시 imageFile 주 렌더러 보존 실패: {img_plan}"


@settings(max_examples=200, deadline=None)
@given(
    has_vertex_image=st.booleans(),
    has_native_diagram=st.booleans(),
    has_image_file=st.booleans(),
    has_slide_bg=st.booleans(),
    role=st.sampled_from(["cover", "section", "structural", "content", "visual"]),
    html_enabled=st.booleans(),
)
def test_pres3_loss_zero_invariant_preserved(has_vertex_image, has_native_diagram,
                                             has_image_file, has_slide_bg, role,
                                             html_enabled):
    """모든 미디어 상태에서 손실-0 불변식 보존: has_vertex_image ⇒ vertex_slot != 'none'.

    **Validates: Requirements 3.2**
    """
    plan = srv._select_render_plan(
        has_vertex_image=has_vertex_image, has_native_diagram=has_native_diagram,
        has_image_file=has_image_file, has_slide_bg=has_slide_bg,
        role=role, html_enabled=html_enabled)
    assert plan["primary"] in ("HTML", "NATIVE_SHAPES", "VERTEX_IMAGE"), plan
    if has_vertex_image:
        assert plan["vertex_slot"] != "none", (
            f"손실-0 위반 — has_vertex_image=True인데 vertex_slot='none': {plan}")
    else:
        assert plan["vertex_slot"] == "none", (
            f"has_vertex_image=False 면 vertex_slot 은 'none' 이어야 함: {plan}")


# ==========================================================================
# PRES-4 (Req 3.4) — 풀블리드 없는 네이티브 텍스트 슬라이드 레이아웃/여백 보존
# ==========================================================================
@settings(max_examples=8, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES))
def test_pres4_plain_native_text_layout_preserved(slide):
    """풀블리드 배경이 없는 비구조 텍스트 슬라이드는 기존 레이아웃·여백을 유지한다:
    풀블리드 PICTURE 없음, 본문 텍스트 보존, 텍스트 도형이 경계 안 + 좌측 여백 유지."""
    title, bullets = slide

    # 전제: 비구조형(다이어그램 아님) → 네이티브 텍스트 경로. Vertex 비활성 → 이미지 없음.
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind == "", f"전제 실패: {title!r}는 비구조형이어야 함 (kind={kind!r})"
    assert not is_bug_condition_scenario({
        "hasBakedFullbleed": False, "hasNativeOverlap": False, "hasOffSlide": False})

    fake, result = _run_html_off([{"title": title, "bullets": list(bullets)}],
                                 vertex_enabled=False)
    slides = _slides(result["absPath"])
    content = slides[1]

    # 풀블리드 배경 없음 + 래스터 이미지 없음(네이티브 텍스트 경로).
    assert _fullbleed_pictures(content) == [], "비구조 텍스트 슬라이드에 풀블리드 배경이 없어야 함"
    assert _pptx_media_bytes(result["absPath"]) == [], (
        "Vertex 비활성 비구조 텍스트 슬라이드에 래스터가 임베드되면 안 됨")
    # 본문 불릿 텍스트 보존.
    texts = _all_shape_text(content)
    hits = [b for b in bullets if b in texts]
    assert hits, f"본문 불릿 텍스트가 보존돼야 함 — bullets={bullets}, texts={texts!r}"
    # 레이아웃/여백 보존: 모든 텍스트 도형이 경계 안 + 본문 텍스트가 좌측 여백 유지(0 flush 아님).
    tshapes = _text_shapes(content)
    assert tshapes, "콘텐츠 슬라이드에 텍스트 도형이 있어야 함"
    for _t, r in tshapes:
        assert _within_bounds(r), f"텍스트 도형이 경계 안이어야 함 — rect={r}"
    body_lefts = [r[0] for t, r in tshapes if any(b in t for b in bullets)]
    assert body_lefts, "본문 텍스트 도형을 찾을 수 있어야 함"
    assert min(body_lefts) >= 0.3 - BOUNDS_EPS, (
        f"본문 텍스트는 좌측 여백을 유지해야 함(edge flush 아님) — left={min(body_lefts):.3f}")


# ==========================================================================
# PRES-5 (Req 3.5) — 비결함 랜덤 덱 보존(슬라이드 수 / 손실-0 / 경계 안)
# ==========================================================================
def _deck_strategy():
    """구조형/콘텐츠/비주얼이 섞인 무작위 비결함 덱(1~4 섹션)."""
    one = st.one_of(
        st.sampled_from(_STRUCTURAL_SLIDES).map(lambda s: {"title": s[0], "bullets": list(s[1])}),
        st.sampled_from(_CONTENT_SLIDES).map(lambda s: {"title": s[0], "bullets": list(s[1])}),
        st.sampled_from(_VISUAL_SLIDES).map(dict),
    )
    return st.lists(one, min_size=1, max_size=4)


@settings(max_examples=12, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(deck=_deck_strategy())
def test_pres5_random_nonbug_deck_preserved(deck):
    """비결함 랜덤 덱: 슬라이드 수(표지+N) 보존, 손실-0(unused==0), 본문 슬라이드 경계 안."""
    assert not is_bug_condition_scenario({
        "hasBakedFullbleed": False, "hasNativeOverlap": False, "hasOffSlide": False})

    fake, result = _run_html_off([dict(s) for s in deck], vertex_enabled=True)
    slides = _slides(result["absPath"])

    # 슬라이드 수 = 표지 + 섹션 수.
    assert len(slides) == 1 + len(deck), (
        f"슬라이드 수 기대 {1 + len(deck)}(표지+{len(deck)}) — 실제 {len(slides)}")
    # 손실-0: 생성된 모든 Vertex 이미지 임베드.
    generated, embedded, unused = _vertex_usage(fake, result["absPath"])
    assert unused == 0, (
        f"손실-0 위반(랜덤 덱): generated={generated}, embedded={embedded}, unused={unused}")
    # 본문 슬라이드(표지 제외) 경계 안.
    for idx, sl in enumerate(slides[1:], start=1):
        off = _off_slide_shapes(sl)
        assert off == [], f"본문 슬라이드 {idx} 에 경계 밖 도형이 없어야 함 — {off}"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
