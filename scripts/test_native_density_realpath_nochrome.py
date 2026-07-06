# Feature: pptx-native-density-render, task21 수정 C: 실경로 로컬 audit(Chrome/게이트웨이 없이)
#
# 목적: server.py 의 실제 생성 진입점 `_tool_generate_pptx` 를 기본 환경
# (AE_NATIVE_LAYOUT_RENDER 미설정 = 기본 활성)에서 헤르메틱하게 구동하여, 다음을
# 입증한다.
#   1) 본문 콘텐츠 슬라이드는 Chrome HTML→PNG 베이크 경로를 *타지 않는다* —
#      `_generate_html_slide_for_section` 이 단 한 번도 호출되지 않음(호출되면 통짜
#      베이크 경로 진입으로 간주해 테스트 실패).
#   2) `_llm_pick_slide_layout`(게이트웨이) 이 (a) 정상 레이아웃을 반환하는 경우와
#      (b) 빈 결과({})를 반환하는 경우, (c) 예외를 던지는 경우(게이트웨이 미가용/
#      타임아웃 시뮬레이션) 모두, 실제 `_tool_generate_pptx` 산출물이
#      `audit_native_density` (a)~(h) 를 전수 통과한다(표지+본문 편집가능·풀블리드
#      통짜 PNG 0·겹침<10%·제목 1회·밀도·베이크 미검출·z-order·스타일).
#   3) 특히 "게이트웨이 빈 결과/예외" 케이스에서도 통짜가 아니라 편집가능 네이티브가
#      나온다(실제 환경 강건성의 핵심 — 근본 원인 해소 입증).
#   4) 수정 B: Chrome 이 가용(_render_html_slide_to_png ok=True)이어도 콘텐츠 구운
#      HTML 표지 PNG 를 풀블리드로 채택하지 않는다(표지 통짜 배경 0, 편집텍스트≥1).
#
# 헤르메틱 규약(네트워크 0·Chrome 0):
#   - 게이트웨이(_get_gw)/LLM(_llm_pick_slide_layout)/Chrome(_find_local_chrome/
#     _render_html_slide_to_png/_generate_html_slide_for_section)/Vertex 전부 mock.
#   - _find_local_chrome 는 truthy mock → _html_enabled=True(프로덕션 동일 마스터
#     게이트). 실제 Chrome 은 띄우지 않는다(렌더 함수 mock).
#
# 실행: ./venv/bin/python -m pytest scripts/test_native_density_realpath_nochrome.py -p no:cacheprovider -q

import os
import sys
import json
import asyncio
import tempfile
import shutil
from unittest import mock

import pytest

_REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _REPO not in sys.path:
    sys.path.insert(0, _REPO)
_SCRIPTS = os.path.join(_REPO, "scripts")
if _SCRIPTS not in sys.path:
    sys.path.insert(0, _SCRIPTS)

import ai_engine.server as srv  # noqa: E402
import audit_pptx_native_density as auditor  # noqa: E402


# ---------------------------------------------------------------------------
# 입력 덱 — 표지 + 본문(콘텐츠) 혼합. 모든 본문 슬라이드가 콘텐츠 텍스트(불릿)를 가짐.
# ---------------------------------------------------------------------------

def _tool_input():
    return {
        "title": "실경로 네이티브 렌더 강건성 검증 덱",
        "highQuality": True,
        "slides": [
            {"title": "프로젝트 개요", "bullets": [
                "목표: 편집 가능한 고밀도 슬라이드 생성",
                "범위: 표지·본문 전 슬라이드",
                "원칙: 통짜 이미지 금지, 네이티브 도형 우선",
                "품질: 겹침 10% 미만, 제목 1회",
            ]},
            {"title": "핵심 지표", "bullets": [
                "처리량 42% 증가",
                "가동률 99.9% 달성",
                "활성 사용자 3.2만",
                "인프라 비용 18% 절감",
            ]},
            {"title": "진행 상태", "bullets": [
                "요구사항 정의 완료",
                "아키텍처 설계 완료",
                "통합 테스트 진행 중",
            ]},
            {"title": "추진 일정", "bullets": [
                "1단계 착수와 킥오프",
                "2단계 핵심 구현",
                "3단계 품질 검증",
                "4단계 단계적 배포",
            ]},
            {"title": "주요 기능", "bullets": [
                "편집 가능한 네이티브 도형",
                "젠스파크급 고밀도 레이아웃",
                "셰이프 겹침 0 보정",
            ]},
        ],
    }


# ---------------------------------------------------------------------------
# 게이트웨이 픽 mock 3종
# ---------------------------------------------------------------------------

def _valid_pick_for(heading: str) -> dict:
    h = heading or ""
    if "지표" in h:
        return {"layout": "kpi_summary", "data": {
            "title": heading, "subtitle": "핵심 성과 지표",
            "features": [
                {"title": "처리량 +42%", "description": "월간 처리량이 크게 늘었습니다."},
                {"title": "가동률 99.9%", "description": "무중단 운영으로 SLA 상회."},
                {"title": "사용자 3.2만", "description": "활성 사용자가 꾸준히 증가."},
                {"title": "비용 -18%", "description": "리소스 최적화로 비용 절감."},
            ]}}
    if "상태" in h:
        return {"layout": "status_table", "data": {
            "title": heading, "left_label": "완료", "left_items": ["요구정의", "설계", "구현"],
            "right_label": "진행", "right_items": ["통합테스트", "튜닝", "검증"]}}
    if "일정" in h:
        return {"layout": "process_flow", "data": {
            "title": heading, "steps": [
                {"label": "1", "title": "착수", "description": "킥오프"},
                {"label": "2", "title": "구현", "description": "반복 개발"},
                {"label": "3", "title": "검증", "description": "품질 게이트"},
                {"label": "4", "title": "배포", "description": "단계적 롤아웃"},
            ]}}
    return {"layout": "two_column", "data": {
        "title": heading, "left_content": "핵심 목표\n일정 준수", "right_content": "리스크 관리\n지속 개선"}}


class _PickNormal:
    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, heading, body, doc="", role="", bullet_count=0):
        self.calls.append(heading)
        return _valid_pick_for(heading)


class _PickEmpty:
    """게이트웨이 빈 결과 시뮬레이션 — layout/data 없음(빈 dict)."""

    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, heading, body, doc="", role="", bullet_count=0):
        self.calls.append(heading)
        return {}


class _PickRaises:
    """게이트웨이 미가용/타임아웃 시뮬레이션 — 예외 발생."""

    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, heading, body, doc="", role="", bullet_count=0):
        self.calls.append(heading)
        raise RuntimeError("gateway unavailable (simulated timeout)")


# ---------------------------------------------------------------------------
# Chrome/HTML mock — 베이크 경로 진입 탐지 스파이 포함
# ---------------------------------------------------------------------------

class _HtmlBakeSpy:
    """_generate_html_slide_for_section 대체 — 호출되면 카운트(=베이크 경로 진입).

    수정 A 후 본문 콘텐츠는 네이티브로 라우팅되어 이 함수가 호출되면 안 된다.
    호출 시 카운트가 증가하고, 테스트가 calls==0 을 단언해 통짜 경로 진입을 감지한다.
    """

    def __init__(self):
        self.calls = 0

    async def __call__(self, *args, **kwargs):
        self.calls += 1
        return ""


async def _render_html_no_chrome(html, output_path, width=1920, height=1080, timeout=30):
    # Chrome 미실행 — 어떤 표지 HTML 도 렌더되지 않는다.
    return {"ok": False, "reason": "mocked-no-chrome"}


def _make_render_html_ok():
    """Chrome 가용 시뮬레이션 — output_path 에 실제 PNG 를 써서 ok=True 반환.

    수정 B 검증용: Chrome 이 표지 HTML PNG 를 성공적으로 렌더해도, 서버는 그 콘텐츠
    구운 PNG 를 풀블리드로 채택하지 않아야 한다.
    """
    async def _render_html_ok(html, output_path, width=1920, height=1080, timeout=30):
        try:
            os.makedirs(os.path.dirname(output_path), exist_ok=True)
            from PIL import Image
            Image.new("RGB", (64, 36), (30, 30, 30)).save(output_path)
        except Exception:
            # PIL 부재 시 최소 바이트라도 기록(존재만 하면 os.path.isfile True).
            with open(output_path, "wb") as f:
                f.write(b"\x89PNG\r\n\x1a\n")
        return {"ok": True}
    return _render_html_ok


def _run(tool_input, project_path, picker, *, render_html=None, env=None):
    """헤르메틱 컨텍스트에서 _tool_generate_pptx 구동 → (result_json, picker, bake_spy)."""
    render_html = render_html or _render_html_no_chrome
    bake_spy = _HtmlBakeSpy()
    env = dict(env or {})
    # 기본 환경: HTML 마스터 게이트 ON, Vertex OFF, 네이티브 라우팅 기본값(미설정).
    env.setdefault("AE_ENABLE_HTML_SLIDES", "1")
    env.setdefault("AE_PREFER_VERTEX_IMAGE", "0")
    env.setdefault("AE_ENABLE_VERTEX_IMAGE", "0")
    env.setdefault("AE_NATIVE_LAYOUT_RENDER", None)  # 미설정 → 기본 활성

    old = {}
    for k, v in env.items():
        old[k] = os.environ.get(k)
        if v is None:
            os.environ.pop(k, None)
        else:
            os.environ[k] = v
    try:
        with mock.patch.object(srv, "_find_local_chrome", lambda: "/fake/chrome"), \
             mock.patch.object(srv, "_call_bridge", lambda ep, payload, timeout=30.0: None), \
             mock.patch.object(srv, "_get_gw", lambda *a, **k: object()), \
             mock.patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "fake-model"), \
             mock.patch.object(srv, "_llm_pick_slide_layout", picker), \
             mock.patch.object(srv, "_render_html_slide_to_png", render_html), \
             mock.patch.object(srv, "_generate_html_slide_for_section", bake_spy):
            out = asyncio.run(srv._tool_generate_pptx(tool_input, project_path))
    finally:
        for k, v in old.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v
    return json.loads(out), picker, bake_spy


def _pptx_path(result, project_path):
    assert "path" in result, f"생성 실패: {result}"
    rel = result["path"]
    root = srv._resolve_local_root(project_path)
    cand = rel if os.path.isabs(rel) else os.path.join(root, rel)
    assert os.path.isfile(cand), f".pptx 부재: {cand}"
    return cand


def _fullbleed_pictures(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    EMU = 914400.0
    n = 0
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            w = (sh.width or 0) / EMU
            h = (sh.height or 0) / EMU
            if w >= 13.0 and h >= 7.2:
                n += 1
        except Exception:
            continue
    return n


def _editable_count(slide):
    n = 0
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                n += 1
        except Exception:
            pass
    return n


def _audit_fail_msg(report):
    return "\n".join(
        f"  - slide {f.get('slide')}: {f.get('check')} shapes={f.get('shapes')} signal={f.get('signal')}"
        for f in report.failures
    )


# ---------------------------------------------------------------------------
# 테스트
# ---------------------------------------------------------------------------

def _assert_audit_and_no_bake(picker_cls):
    tmp = tempfile.mkdtemp(prefix="ae_nochrome_")
    try:
        result, pick, bake_spy = _run(_tool_input(), tmp, picker_cls())
        path = _pptx_path(result, tmp)

        # (1) 통짜 베이크 경로 미진입 — _generate_html_slide_for_section 미호출.
        assert bake_spy.calls == 0, (
            f"본문 HTML 베이크 경로 진입 감지: _generate_html_slide_for_section "
            f"{bake_spy.calls}회 호출됨(통짜 이미지 경로)")

        # (2) 산출물 audit (a)~(h) 전수 통과.
        report = auditor.audit_native_density(path)
        assert report.passed, "실경로 audit 미통과:\n" + _audit_fail_msg(report)
        assert report.failures == []

        # (3) 전 슬라이드 편집가능 텍스트 ≥1 + 풀블리드 통짜 배경 0.
        from pptx import Presentation
        prs = Presentation(path)
        for si, slide in enumerate(prs.slides):
            assert _editable_count(slide) >= 1, f"슬라이드 {si + 1} 편집가능 텍스트 0(통짜 의심)"
            assert _fullbleed_pictures(slide) == 0, f"슬라이드 {si + 1} 풀블리드 통짜 배경 채택됨"
        return pick
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def test_gateway_normal_pick_realpath_audit_passes_no_bake():
    """게이트웨이 정상 픽 — audit (a)~(h) 통과 + 베이크 미진입."""
    pick = _assert_audit_and_no_bake(_PickNormal)
    assert pick.calls, "네이티브 라우팅이 _llm_pick_slide_layout 을 호출하지 않음"


def test_gateway_empty_result_realpath_audit_passes_no_bake():
    """게이트웨이 빈 결과({}) — 통짜가 아니라 편집가능 네이티브로 귀결하고 audit 통과.

    이 케이스가 실제 환경 강건성의 핵심(근본 원인 해소): 게이트웨이가 아무 레이아웃도
    주지 못해도 합성 레이아웃으로 반드시 편집가능 네이티브 렌더."""
    pick = _assert_audit_and_no_bake(_PickEmpty)
    assert pick.calls, "빈결과 케이스에서도 픽은 시도되어야 함"


def test_gateway_raises_realpath_audit_passes_no_bake():
    """게이트웨이 예외(미가용/타임아웃) — 예외 흡수 후에도 편집가능 네이티브·audit 통과."""
    pick = _assert_audit_and_no_bake(_PickRaises)
    assert pick.calls, "예외 케이스에서도 픽은 시도되어야 함"


def test_cover_html_not_baked_even_when_chrome_available():
    """수정 B — Chrome 가용(render_html ok=True)이어도 콘텐츠 구운 HTML 표지 PNG 를
    풀블리드로 채택하지 않는다(표지 통짜 배경 0, 편집 텍스트 ≥1)."""
    tmp = tempfile.mkdtemp(prefix="ae_nochrome_cover_")
    try:
        result, pick, bake_spy = _run(
            _tool_input(), tmp, _PickNormal(), render_html=_make_render_html_ok())
        path = _pptx_path(result, tmp)

        from pptx import Presentation
        prs = Presentation(path)
        cover = list(prs.slides)[0]
        # 콘텐츠 구운 풀블리드 표지 배경 미채택.
        assert _fullbleed_pictures(cover) == 0, "표지에 콘텐츠 구운 HTML 풀블리드 배경이 채택됨(수정 B 위반)"
        # 표지 편집 가능 텍스트 존재.
        assert _editable_count(cover) >= 1, "표지 편집 가능 텍스트 0(통짜 의심)"
        # 본문도 여전히 베이크 미진입.
        assert bake_spy.calls == 0, f"본문 베이크 경로 진입: {bake_spy.calls}회"

        report = auditor.audit_native_density(path)
        assert report.passed, "Chrome 가용 케이스 audit 미통과:\n" + _audit_fail_msg(report)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
