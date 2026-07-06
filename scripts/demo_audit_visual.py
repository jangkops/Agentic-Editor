"""before.pptx / after.pptx 를 audit 도구로 검증해 수치를 출력(육안 확인 보조)."""
import io, os, sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)
import audit_pptx_zorder_break as azb

SW, SH, EPS = 13.333, 7.5, 0.05


def _px(sh):
    try:
        with Image.open(io.BytesIO(sh.image.blob)) as im:
            return im.size
    except Exception:
        return (0, 0)


def report(path):
    print(f"\n=== {os.path.basename(path)} ===")
    for i, slide in enumerate(Presentation(path).slides, start=1):
        fb = off = small_large = 0
        for sh in slide.shapes:
            try:
                if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue
            r = azb._rect(sh)
            if r is None:
                continue
            if azb._fullbleed(r):
                fb += 1
            if r[0] < -EPS or r[1] < -EPS or r[0]+r[2] > SW+EPS or r[1]+r[3] > SH+EPS:
                off += 1
            pw, ph = _px(sh)
            if (pw >= 1024 or ph >= 1024) and r[2] <= 0.5 and r[3] <= 0.5:
                small_large += 1
        tag = {1: "D1", 2: "D2", 3: "D3"}.get(i, f"S{i}")
        print(f"  슬라이드{i}({tag}): 풀블리드={fb}  슬라이드밖={off}  "
              f"소형슬롯_대형이미지={small_large}")


for name in ("before.pptx", "after.pptx"):
    report(os.path.join(_ROOT, ".generated", "_visual_proof", name))
