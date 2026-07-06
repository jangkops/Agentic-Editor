"""Property-based test — pptx-native-density-render (작업 15.1).

이 파일은 산출물_검증기 ``scripts/audit_pptx_native_density.audit_native_density`` 의
**합격/불합격 결정성과 실패 보고 완전성**(Correctness Property 15)을 Hypothesis 로
검증한다. Property 15 는 **단일** property-based test 로 구현하며
``@settings(max_examples=...)`` 로 충분한 반복을 보장한다.

핵심 불변식(Property 15):
  - (a)~(g) 검사가 전부 통과할 때에만 ``passed=True`` 이고, 하나라도 실패하면
    ``passed=False`` 이며 ``failures`` 가 비어있지 않다. 즉
    ``report.passed == (len(report.failures) == 0)`` (결정성).
  - 각 failure 는 ``{check, slide(1-base ≥1), shapes(list), signal}`` 키를 모두
    보유한다(보고 완전성).
  - 결함을 주입한 슬라이드(경계 밖 도형 / 겹치는 텍스트 2개 / 동일 제목 2개)에는
    대응 검사 항목(c_bounds / b_overlap / d_title_once)이 해당 슬라이드 번호로
    failures 에 나타난다.

실행 규약(design.md §Testing Strategy):
  - **실제 .pptx I/O 포함**: 깨끗한 네이티브 덱(`render_native_layout`)과 결함 주입
    덱을 python-pptx 로 실제 .pptx 파일에 저장한 뒤 audit 한다. 임시 파일은
    OS 임시 경로에 쓰고 매 예제마다 정리한다.
  - **네트워크 0 / Chrome 0**: 게이트웨이/Vertex 실호출 없음, Vertex 옵트인은 OFF 로
    강제(`AE_ENABLE_VERTEX_IMAGE` != "1") → 장식 이미지 미생성, 콘텐츠 네이티브만.
    audit 는 python-pptx 만으로 동작(Chrome 불필요).
  - 실제 .pptx 저장/열기 비용이 있어 deadline=None, max_examples 는 60(≥50 유지).
    20초+ hang 시 max_examples 를 낮춘다.

실행:
  ./venv/bin/python -m pytest scripts/test_native_density_audit_pbt.py -p no:cacheprovider -q
"""

from __future__ import annotations

import os
import sys
import string
import tempfile
from collections import defaultdict

# --- Vertex 옵트인 OFF 강제 (네트워크 0, 장식 이미지 미생성) ---
os.environ["AE_ENABLE_VERTEX_IMAGE"] = "0"

# --- import 경로: repo 루트(ai_engine 패키지) + scripts 디렉터리(audit 모듈) ---
_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.dirname(_THIS_DIR)
for _p in (_REPO_ROOT, _THIS_DIR):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest
from hypothesis import given, settings, strategies as st, HealthCheck

from pptx import Presentation
from pptx.util import Inches, Pt

from ai_engine.native_layout_renderer import (
    render_native_layout,
    NATIVE_LAYOUT_REGISTRY,
)
from ai_engine.slide_templates import design_tokens_for_profile

# 대상 산출물_검증기 + 검사 항목 식별자(결함 주입 → 대응 check 검증용).
from audit_pptx_native_density import (
    audit_native_density,
    AuditReport,
    CHECK_OVERLAP,      # "b_overlap"
    CHECK_BOUNDS,       # "c_bounds"
    CHECK_TITLE_ONCE,   # "d_title_once"
)

_TOKENS = design_tokens_for_profile(None)
_KNOWN_LAYOUTS = sorted(NATIVE_LAYOUT_REGISTRY.keys())

_COMMON_SETTINGS = settings(
    max_examples=60,          # ≥50 유지(실제 .pptx I/O 비용 반영, deadline=None)
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow, HealthCheck.data_too_large,
                           HealthCheck.filter_too_much],
)

# failures 원소가 반드시 보유해야 하는 키(보고 완전성).
_REQUIRED_FAILURE_KEYS = {"check", "slide", "shapes", "signal"}


# ===========================================================================
# 텍스트/데이터 생성기 (한글·특수문자·빈선택·장문 포함)
# ===========================================================================

# python-pptx XML 에 안전한 문자만(제어문자/서로게이트 배제).
_SAFE_ALPHABET = (
    "가나다라마바사아자차카타파하거너더러머버서어전정종주중지"
    + string.ascii_letters
    + string.digits
    + " .,!?·—…:/@#%&*()-_=+[]{}"
)


def _ne_text(min_size: int = 1, max_size: int = 40):
    return (st.text(alphabet=_SAFE_ALPHABET, min_size=min_size, max_size=max_size)
            .map(lambda s: s.strip())
            .filter(lambda s: len(s) > 0))


_long_text = (st.text(alphabet=_SAFE_ALPHABET, min_size=120, max_size=220)
              .map(lambda s: s.strip()).filter(lambda s: len(s) > 0))

_content_text = st.one_of(_ne_text(1, 40), _ne_text(1, 80), _long_text)
_optional_text = st.one_of(st.just(""), st.just("   "), _ne_text(1, 30))


def _cards(min_size=2, max_size=3):
    return st.lists(st.fixed_dictionaries({
        "icon": st.just("★"),
        "title": _ne_text(1, 28),
        "description": _optional_text,
    }), min_size=min_size, max_size=max_size)


def _steps(min_size=1, max_size=3):
    return st.lists(st.fixed_dictionaries({
        "label": _ne_text(1, 12),
        "title": _ne_text(1, 28),
        "description": _optional_text,
    }), min_size=min_size, max_size=max_size)


def _items(min_size=1, max_size=3):
    return st.lists(_ne_text(1, 36), min_size=min_size, max_size=max_size)


def _layers(min_size=1, max_size=3):
    return st.lists(st.fixed_dictionaries({
        "name": _ne_text(1, 26),
        "description": _optional_text,
        "items": st.lists(_ne_text(1, 18), min_size=0, max_size=2),
    }), min_size=min_size, max_size=max_size)


def _layout_data_strategy(layout: str):
    """알려진_레이아웃별 유효 data dict 생성기(필수 충족 + 한글/특수/장문/빈선택)."""
    if layout == "cover":
        return st.fixed_dictionaries(
            {"title": _content_text},
            optional={"subtitle": _optional_text, "eyebrow": _optional_text,
                      "footer": _optional_text})
    if layout == "section_divider":
        return st.fixed_dictionaries(
            {"title": _content_text},
            optional={"section_number": st.integers(1, 9), "description": _optional_text})
    if layout == "two_column":
        return st.fixed_dictionaries(
            {"title": _content_text,
             "left_content": st.one_of(_content_text, _items(1, 3)),
             "right_content": st.one_of(_content_text, _items(1, 3))},
            optional={"subtitle": _optional_text})
    if layout == "feature_grid":
        return st.fixed_dictionaries(
            {"title": _content_text, "features": _cards(2, 3)},
            optional={"subtitle": _optional_text})
    if layout == "timeline":
        return st.fixed_dictionaries(
            {"title": _content_text, "steps": _steps(1, 3)},
            optional={"subtitle": _optional_text})
    if layout == "comparison":
        return st.fixed_dictionaries(
            {"title": _content_text,
             "left_label": _ne_text(1, 18), "left_items": _items(1, 3),
             "right_label": _ne_text(1, 18), "right_items": _items(1, 3)},
            optional={"subtitle": _optional_text})
    if layout == "architecture":
        return st.fixed_dictionaries(
            {"title": _content_text, "layers": _layers(1, 3)},
            optional={"subtitle": _optional_text})
    raise AssertionError(f"unknown layout {layout}")


@st.composite
def _slide_spec(draw):
    """덱의 한 슬라이드 명세 — 깨끗한 네이티브 렌더 또는 결함 주입 3종 중 하나."""
    kind = draw(st.sampled_from(
        ["clean", "defect_bounds", "defect_overlap", "defect_title"]))
    if kind == "clean":
        layout = draw(st.sampled_from(_KNOWN_LAYOUTS))
        data = draw(_layout_data_strategy(layout))
        return {"kind": "clean", "layout": layout, "data": data}
    # 결함 슬라이드는 결정적으로 구성(추가 데이터 불필요).
    if kind == "defect_title":
        return {"kind": kind, "title": draw(_ne_text(1, 24))}
    return {"kind": kind}


# ===========================================================================
# 슬라이드 빌더 (python-pptx 직접 — 결함 주입은 결정적 좌표/텍스트)
# ===========================================================================


def _add_textbox(slide, l, t, w, h, text, font_pt=None):
    """편집가능 텍스트박스 1개 추가(필요 시 명시 폰트 크기)."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = text
    if font_pt is not None:
        for p in tf.paragraphs:
            if not p.runs:
                continue
            for r in p.runs:
                r.font.size = Pt(font_pt)
    return tb


def _build_slide(prs, spec):
    """명세대로 빈 슬라이드에 콘텐츠/결함을 구성한다."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    kind = spec["kind"]

    if kind == "clean":
        render_native_layout(slide, prs, spec["layout"], spec["data"], _TOKENS)
        return slide

    if kind == "defect_bounds":
        # 정상 콘텐츠(경계 안) 1개 + 경계를 명백히 벗어난 도형 1개 → c_bounds 보장.
        _add_textbox(slide, 1.0, 0.6, 6.0, 0.8, "경계 안 제목", font_pt=28)
        # left=14in → 슬라이드 폭 13.333in + eps 0.05 초과 → within_bounds 실패.
        _add_textbox(slide, 14.0, 1.0, 2.5, 1.0, "out of bounds", font_pt=14)
        return slide

    if kind == "defect_overlap":
        # 텍스트 보유 2개가 크게 겹침(겹침률 ≈ 56% ≥ 10%) → b_overlap 보장.
        _add_textbox(slide, 2.0, 2.0, 4.0, 2.0, "겹치는 텍스트 A", font_pt=18)
        _add_textbox(slide, 3.0, 2.5, 4.0, 2.0, "겹치는 텍스트 B", font_pt=18)
        return slide

    if kind == "defect_title":
        # 동일 텍스트·동일 최대폰트 2개(비겹침) → 제목 셰이프 2개 → d_title_once 보장.
        t = spec["title"]
        _add_textbox(slide, 1.0, 0.5, 5.0, 0.7, t, font_pt=32)
        _add_textbox(slide, 1.0, 2.2, 5.0, 0.7, t, font_pt=32)
        return slide

    raise AssertionError(f"unknown spec kind {kind}")


def _make_deck(specs):
    """16:9(13.333×7.5) Presentation 에 명세대로 슬라이드를 채워 반환."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for spec in specs:
        _build_slide(prs, spec)
    return prs


# ===========================================================================
# Property 15 (작업 15.1) — 단일 PBT
# ===========================================================================
# Feature: pptx-native-density-render, Property 15: 산출물 검증기의 합격/불합격 결정성과 보고
@_COMMON_SETTINGS
@given(specs=st.lists(_slide_spec(), min_size=1, max_size=3))
def test_property15_audit_determinism_and_report_completeness(specs):
    """For any 생성된 .pptx 에 대해 audit_native_density 는:

      1. (a)~(g) 전부 통과 시에만 passed=True — 즉 passed == (failures 비어있음).
      2. 각 failure 는 {check, slide(1-base ≥1), shapes(list), signal} 키를 보유.
      3. 결함 주입 슬라이드(경계 밖/겹침 2개/동일 제목 2개)에는 대응 검사 항목이
         해당 슬라이드 번호로 failures 에 나타난다.

    **Validates: Requirements 8.1, 8.2, 8.3, 8.4, 8.6, 8.8, 8.10**
    """
    prs = _make_deck(specs)
    n_slides = len(specs)

    fd, path = tempfile.mkstemp(suffix=".generated.pptx", prefix="pbt_audit_native_")
    os.close(fd)
    try:
        prs.save(path)
        report = audit_native_density(path)
    finally:
        try:
            os.remove(path)
        except OSError:
            pass

    # --- 타입 ---
    assert isinstance(report, AuditReport)
    assert isinstance(report.passed, bool)
    assert isinstance(report.failures, list)

    # --- 불변식 1: 결정성 — passed == (failures 비어있음) ---
    assert report.passed == (len(report.failures) == 0)

    # --- 불변식 2: 각 failure 보고 완전성 ---
    for f in report.failures:
        assert isinstance(f, dict)
        assert _REQUIRED_FAILURE_KEYS.issubset(f.keys()), \
            f"failure 키 누락: {sorted(f.keys())}"
        # slide 는 1-base 정수이며 덱 범위 안.
        assert isinstance(f["slide"], int)
        assert 1 <= f["slide"] <= n_slides, f"slide 범위 위반: {f['slide']} (n={n_slides})"
        # shapes 는 리스트.
        assert isinstance(f["shapes"], list)
        # signal 은 비어있지 않은 문자열.
        assert isinstance(f["signal"], str) and f["signal"]
        # check 는 비어있지 않은 문자열.
        assert isinstance(f["check"], str) and f["check"]

    # --- 불변식 3: 결함 주입 → 대응 검사 항목이 해당 슬라이드에 나타남 ---
    checks_by_slide = defaultdict(set)
    for f in report.failures:
        checks_by_slide[f["slide"]].add(f["check"])

    for idx, spec in enumerate(specs, start=1):
        kind = spec["kind"]
        if kind == "defect_bounds":
            assert CHECK_BOUNDS in checks_by_slide[idx], \
                f"슬라이드 {idx}: 경계 밖 도형인데 {CHECK_BOUNDS} 미보고 ({checks_by_slide[idx]})"
        elif kind == "defect_overlap":
            assert CHECK_OVERLAP in checks_by_slide[idx], \
                f"슬라이드 {idx}: 겹침 텍스트인데 {CHECK_OVERLAP} 미보고 ({checks_by_slide[idx]})"
        elif kind == "defect_title":
            assert CHECK_TITLE_ONCE in checks_by_slide[idx], \
                f"슬라이드 {idx}: 동일 제목 2개인데 {CHECK_TITLE_ONCE} 미보고 ({checks_by_slide[idx]})"

    # --- 결함 주입 슬라이드가 하나라도 있으면 반드시 불합격(결정성 교차 확인) ---
    if any(s["kind"] != "clean" for s in specs):
        assert report.passed is False
        assert len(report.failures) >= 1
