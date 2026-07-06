"""Regression — 템플릿은 *스타일*로만 상속하고 내용은 새로 생성한다.

사용자 보고: "템플릿을 활용하여 작성"을 요청하면, 템플릿의 기존 샘플 슬라이드 위에
새 슬라이드가 *추가*되어 기존 내용과 새 내용이 섞이는(꼬이는) 버그.

근본 원인: `_tool_generate_pptx`가 templatePath의 .pptx를 열면 그 안의 샘플 슬라이드가
그대로 남고, cover/content 슬라이드가 뒤에 append됨.

수정: 템플릿을 연 직후 `_clear_all_slides(prs)`로 기존 슬라이드를 모두 제거(마스터/
레이아웃/테마는 보존) → 스타일만 상속하고 내용은 새로 생성.

Correctness properties:
  P1. _clear_all_slides는 슬라이드만 제거하고 레이아웃/마스터는 보존한다.
  P2. 템플릿으로 생성한 PPTX의 슬라이드 수 == 1(cover) + len(slides). (템플릿 샘플 미포함)
  P3. 템플릿 샘플 슬라이드의 내용(마커)은 결과물에 남지 않는다.

실행: pytest scripts/test_pptx_template_clears_slides.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_clear_all_slides_keeps_layouts():
    """P1 — 슬라이드만 제거, 레이아웃/마스터 보존(이후 add_slide 가능)."""
    from pptx import Presentation
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides.add_slide(prs.slide_layouts[1])
    assert len(prs.slides._sldIdLst) == 2
    n_layouts = len(prs.slide_layouts)

    removed = server._clear_all_slides(prs)
    assert removed == 2
    assert len(prs.slides._sldIdLst) == 0
    # 레이아웃 보존 → 계속 사용 가능
    assert len(prs.slide_layouts) == n_layouts
    prs.slides.add_slide(prs.slide_layouts[0])
    assert len(prs.slides._sldIdLst) == 1


def test_template_sample_slides_not_mixed(tmp_path):
    """P2+P3 — 템플릿 샘플 슬라이드가 결과물에 섞이지 않는다."""
    from pptx import Presentation

    # 샘플 슬라이드 2개가 든 '템플릿' .pptx 생성
    tpl = Presentation()
    s1 = tpl.slides.add_slide(tpl.slide_layouts[0])
    s1.shapes.title.text = "TEMPLATE_SAMPLE_MARKER_1"
    s2 = tpl.slides.add_slide(tpl.slide_layouts[1])
    s2.shapes.title.text = "TEMPLATE_SAMPLE_MARKER_2"
    tpl_path = tmp_path / "tpl.pptx"
    tpl.save(str(tpl_path))

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    slides = [
        {"title": "새 내용 1", "bullets": ["항목 a"]},
        {"title": "새 내용 2", "bullets": ["항목 b"]},
    ]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "새 덱", "slides": slides, "templatePath": str(tpl_path)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        prs = Presentation(res["absPath"])
        # P2 — cover(1) + content(2) = 3. 템플릿 샘플 2개가 섞였다면 5가 된다.
        assert len(prs.slides._sldIdLst) == 3, (
            f"슬라이드 수 {len(prs.slides._sldIdLst)} — 템플릿 샘플이 섞임(기대 3)")
        # P3 — 템플릿 샘플 마커가 결과물에 없어야 한다
        all_text = []
        for sl in prs.slides:
            for sh in sl.shapes:
                if sh.has_text_frame:
                    all_text.append(sh.text_frame.text)
        joined = "\n".join(all_text)
        assert "TEMPLATE_SAMPLE_MARKER" not in joined, "템플릿 샘플 내용이 결과물에 잔존"
        assert "새 덱" in joined, "새로 생성한 표지 제목이 없음"
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
