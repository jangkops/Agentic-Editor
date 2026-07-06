"""Property-based tests — spec: pptx-native-density-render, Task 16 (Vertex 장식·z-order).

이 파일은 ``ai_engine/native_layout_renderer`` 의 Vertex 장식_비주얼 경로
(``maybe_generate_decorative`` / ``emit_decorative_figure`` /
``maybe_add_decorative_background`` / ``_get_vertex_client`` / ``_send_to_back``)와
산출물_검증기 ``scripts/audit_pptx_native_density.audit_native_density`` ·
``scripts/audit_pptx_baked_text.baked_text_score`` 에 대한 5개 Correctness Property
(12·13·16·17·18)를 각각 단일 Hypothesis PBT 로 검증한다.

각 Property = 단일 테스트, ``@settings(max_examples=100)`` 이상.

**Vertex mock 필수 — 실제 Vertex/네트워크 호출 절대 금지.**
  - ``native_layout_renderer._get_vertex_client`` 를 monkeypatch 하여 가짜 클라이언트
    (``enabled``/``generate``)를 주입한다. 가짜 ``generate`` 는 PIL 또는 1x1 PNG 의
    base64 를 ``{"images":[...]}`` 로 반환한다(네트워크 0).
  - 옵트인 게이트(``AE_ENABLE_VERTEX_IMAGE``)는 테스트 내부에서 켜고(ON="1") 끈다
    (OFF=미설정/"0") — 컨텍스트 매니저로 복원.

실행 환경: python-pptx in-memory + PIL 만 사용한다(Chrome 불필요·네트워크 0).
실행:
  ./venv/bin/python -m pytest scripts/test_native_vertex_decorative_pbt.py -p no:cacheprovider -q
"""
from __future__ import annotations

import io
import os
import sys
import base64
import shutil
import tempfile
import contextlib
from unittest.mock import patch

# repo root + ai_engine + scripts 를 path 에 올려 native_layout_renderer(ai_engine)
# 와 audit_pptx_native_density / audit_pptx_baked_text(scripts) 를 모두 import 가능하게 한다.
_HERE = os.path.dirname(os.path.abspath(__file__))
_REPO = os.path.dirname(_HERE)
_AE = os.path.join(_REPO, "ai_engine")
for _p in (_REPO, _AE, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import native_layout_renderer as nlr  # noqa: E402  (ai_engine on path)
import audit_pptx_native_density as aud  # noqa: E402  (scripts on path)
from audit_pptx_baked_text import baked_text_score  # noqa: E402


# ===========================================================================
# 공용 상수 / 헬퍼
# ===========================================================================

SLIDE_W = 13.333
SLIDE_H = 7.5
EMU_PER_IN = 914400.0

# 베이크 텍스트 판정(audit_pptx_baked_text 와 동일): pct>=6.0 OR lines>=6 → 초과.
_BAKED_PCT = 6.0
_BAKED_LINES = 6

# 풀블리드 배경 후보 이미지 크기(px). 1200 이하라 baked_text_score 가 리사이즈하지
# 않는다(에지 보존). 너무 크면 느리므로 적당히.
_IMG_W = 360
_IMG_H = 260

# 콘텐츠 텍스트 풀(짧고 고유한 문자열 — pptx XML 안전).
_CONTENT_POOL = [
    "데이터 수집", "정제 단계", "분석 파이프라인", "시각화 산출",
    "모델 학습", "검증 절차", "배포 자동화", "모니터링 대시보드",
    "Alpha module", "Beta service", "Gamma layer", "Delta report",
]
_TITLE_TEXT = "프레젠테이션 제목 헤더"   # 콘텐츠 풀과 겹치지 않는 고유 제목
_SUBTITLE_TEXT = "부제목 요약 문장"


def _tokens() -> dict:
    """기존 design_tokens 재사용(신규 토큰 정의 없음)."""
    return nlr.design_tokens_for_profile(None)


def _new_slide():
    """13.333×7.5in 빈(blank) 슬라이드 1개를 가진 Presentation/Slide 생성."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return prs, slide


@contextlib.contextmanager
def _optin_env(value):
    """AE_ENABLE_VERTEX_IMAGE 를 설정/복원. value=None 이면 키 삭제(OFF)."""
    key = "AE_ENABLE_VERTEX_IMAGE"
    old = os.environ.get(key)
    if value is None:
        os.environ.pop(key, None)
    else:
        os.environ[key] = value
    try:
        yield
    finally:
        if old is None:
            os.environ.pop(key, None)
        else:
            os.environ[key] = old


class _FakeVertexClient:
    """Vertex 이미지 클라이언트 mock — 실제 Vertex/네트워크 호출 없음.

    ``generate`` 는 async(실제 모듈과 동일)로, 주어진 PNG 바이트의 base64 를
    ``{"images":[b64]}`` 로 반환한다. ``enabled``/``fail``/``error`` 로 옵트인/실패
    경로를 흉내낸다.
    """

    def __init__(self, png_bytes: bytes = b"", *, enabled: bool = True,
                 fail: bool = False, error: bool = False):
        self._png = png_bytes
        self.enabled = enabled
        self.fail = fail
        self.error = error
        self.calls = 0

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None, **_kw):
        self.calls += 1
        if self.error:
            return {"error": "mock-error (no network)"}
        if self.fail:
            raise RuntimeError("mock-generate-failure (no network)")
        return {"images": [base64.b64encode(self._png).decode("ascii")]}


@contextlib.contextmanager
def _vertex(fake: _FakeVertexClient, *, optin="1"):
    """_get_vertex_client 를 fake 로 patch + 옵트인 env 설정(컨텍스트)."""
    with patch.object(nlr, "_get_vertex_client", lambda **_kw: fake), _optin_env(optin):
        yield


def _clean_decor_tmp():
    """maybe_*_decorative 의 기본 임시 이미지 디렉터리 정리($TMPDIR/ae_vertex_decor)."""
    shutil.rmtree(os.path.join(tempfile.gettempdir(), "ae_vertex_decor"),
                  ignore_errors=True)


# --- 이미지 생성기 (PIL, 네트워크 0) --------------------------------------

def _abstract_png(seed: int) -> bytes:
    """장식_배경(콘텐츠 텍스트 없는 추상 이미지) — 단색/평탄 → 에지 거의 0.

    baked_text_score 판정을 확실히 미초과(텍스트추정행≈0%, 줄≈0)한다.
    """
    img = Image.new("RGB", (_IMG_W, _IMG_H),
                    (30 + seed % 200, 30 + (seed * 3) % 200, 30 + (seed * 7) % 200))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _banded_png(nbands: int, band_h: int) -> bytes:
    """텍스트가 구워진 듯한 이미지 — 다수의 고주파 세로 줄무늬 밴드.

    각 밴드(연속 band_h 행)는 흑/배경 줄무늬로 행당 에지 급변이 매우 많아
    '텍스트 줄'로 추정된다. 밴드 수/높이에 따라 baked_text_score 가 임계를
    넘기도/못 넘기도 한다(경계 ± 포함). 실제 점수는 테스트에서 계산해 분기한다.
    """
    img = Image.new("L", (_IMG_W, _IMG_H), 210)  # 평탄 배경(밝음)
    d = ImageDraw.Draw(img)
    stripe_w = 2
    gap = 4
    y = 5
    placed = 0
    while placed < nbands and y + band_h < _IMG_H:
        x = 0
        while x < _IMG_W:
            d.rectangle([x, y, x + stripe_w - 1, y + band_h - 1], fill=0)
            x += 2 * stripe_w  # 흑 줄무늬 + 동일 폭 배경 갭
        y += band_h + gap
        placed += 1
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


def _tiny_png(seed: int, w: int = 8, h: int = 8) -> bytes:
    """내용 무관한 작은 유효 PNG(장식 슬롯 채움/손실-0 검증용)."""
    img = Image.new("RGB", (max(1, w), max(1, h)),
                    (seed % 256, (seed * 5) % 256, (seed * 11) % 256))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


# --- 슬라이드 검사 헬퍼 ---------------------------------------------------

def _is_picture(sh) -> bool:
    try:
        return sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _pictures(slide):
    """[(z_index, shape)] — 슬라이드의 PICTURE 셰이프(배치=z 순서)."""
    return [(i, sh) for i, sh in enumerate(slide.shapes) if _is_picture(sh)]


def _content_text_shapes(slide):
    """[(z_index, shape)] — 비어있지 않은 텍스트 런을 가진 편집가능 텍스트 셰이프."""
    out = []
    for i, sh in enumerate(slide.shapes):
        try:
            if _is_picture(sh):
                continue
            if getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                out.append((i, sh))
        except Exception:
            continue
    return out


def _all_editable_text(slide) -> str:
    return " \n ".join(sh.text_frame.text for _, sh in _content_text_shapes(slide))


def _rect_in(sh):
    """셰이프 경계 사각형을 인치 (left, top, width, height) 로."""
    return (sh.left / EMU_PER_IN, sh.top / EMU_PER_IN,
            sh.width / EMU_PER_IN, sh.height / EMU_PER_IN)


def _within_bounds(rect, eps: float = 0.05) -> bool:
    l, t, w, h = rect
    return (l >= -eps and t >= -eps
            and l + w <= SLIDE_W + eps and t + h <= SLIDE_H + eps)


def _overlap_ratio(a, b) -> float:
    """min 면적 대비 교집합 비율(0~1)."""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0.0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0.0, min(ay + ah, by + bh) - max(ay, by))
    inter = ix * iy
    amin = min(max(0.0, aw * ah), max(0.0, bw * bh))
    return (inter / amin) if amin > 0 else 0.0


def _save_pptx(prs, tmpdir) -> str:
    path = os.path.join(tmpdir, "deck.pptx")
    prs.save(path)
    return path


def _baked_exceeded(png_bytes: bytes) -> bool:
    """baked_text_score 판정 초과 여부(pct>=6 OR lines>=6)."""
    pct, lines = baked_text_score(Image.open(io.BytesIO(png_bytes)))
    return pct >= _BAKED_PCT or lines >= _BAKED_LINES


def _slide_baked_failures(report, slide_no: int):
    return [f for f in report.failures
            if f.get("check") == aud.CHECK_BAKED_TEXT and f.get("slide") == slide_no]


def _slide_zorder_failures(report, slide_no: int):
    return [f for f in report.failures
            if f.get("check") == aud.CHECK_ZORDER and f.get("slide") == slide_no]


_SUPPRESS = [HealthCheck.too_slow, HealthCheck.data_too_large]


# ===========================================================================
# Property 12 (Task 16.1) — 풀블리드 배경은 장식만 허용된다 (Req 6.1, 6.3)
# ===========================================================================
# Feature: pptx-native-density-render, Property 12: 풀블리드 배경 후보 이미지에 대해,
# baked_text_score 판정(텍스트추정행 비율 < 6% AND 추정 텍스트줄 < 6)을 만족하는
# 장식_배경만 풀블리드로 채택되고, 판정을 초과(비율 >= 6% OR 줄 >= 6)하는 이미지는
# 풀블리드로 채택되지 않으며 해당 슬라이드 콘텐츠는 편집가능_네이티브로 렌더된다.
@settings(max_examples=100, deadline=None, suppress_health_check=_SUPPRESS)
@given(
    kind=st.sampled_from(["abstract", "banded"]),
    nbands=st.integers(min_value=3, max_value=10),
    band_h=st.integers(min_value=3, max_value=5),
    seed=st.integers(min_value=0, max_value=255),
)
def test_p12_fullbleed_background_only_decorative(kind, nbands, band_h, seed):
    """풀블리드 배경 채택은 정확히 baked_text_score 판정에 의해 결정된다.

    Vertex mock 이 후보 이미지를 풀블리드 장식_배경으로 생성하고(손실-0 임베드),
    render_native_layout 이 그 위에 편집가능 콘텐츠를 렌더한다. audit_native_density 의
    (f) 베이크 텍스트 검사는 **판정 초과 이미지에만** 실패를 보고하고, 미초과(추상)
    이미지에는 보고하지 않는다(=풀블리드로 정상 채택). 어느 경우든 콘텐츠는 편집가능
    네이티브로 공존한다(Req 6.3: 초과 시 콘텐츠를 네이티브로 렌더).
    """
    png = _abstract_png(seed) if kind == "abstract" else _banded_png(nbands, band_h)
    expected_baked = _baked_exceeded(png)  # 실제 점수로 기대값 산출(경계 정확).

    fake = _FakeVertexClient(png, enabled=True)
    data = {"title": _TITLE_TEXT, "eyebrow": "EYEBROW",
            "subtitle": _SUBTITLE_TEXT, "footer": "2026"}
    tmp = tempfile.mkdtemp()
    try:
        with _vertex(fake, optin="1"):
            prs, slide = _new_slide()
            res = nlr.render_native_layout(slide, prs, "cover", data, _tokens())
            assert res.ok, "cover 렌더는 성공해야 함"
            path = _save_pptx(prs, tmp)

        # 풀블리드 장식_배경(Vertex mock)이 실제로 채택되었는지(=임베드) 확인.
        pics = _pictures(slide)
        assert len(pics) == 1, f"풀블리드 장식_배경 1개가 임베드돼야 함 (pics={len(pics)})"
        fb_idx, fb_sh = pics[0]
        assert _within_bounds(_rect_in(fb_sh)), "풀블리드 배경은 슬라이드 경계 안이어야 함"
        # 손실-0: 임베드 바이트 == mock 이 반환한 원본 바이트.
        assert fb_sh.image.blob == png, "장식_배경은 손실-0(원본 바이트 동일)으로 임베드돼야 함"

        # Req 6.3: 콘텐츠는 (배경 채택 여부와 무관하게) 편집가능 네이티브로 렌더된다.
        content = _content_text_shapes(slide)
        assert len(content) >= 1, "콘텐츠는 편집가능 네이티브 텍스트로 렌더돼야 함"
        assert _TITLE_TEXT in _all_editable_text(slide), "제목이 편집가능 텍스트로 보존돼야 함"

        # 핵심 불변식: (f) 베이크 텍스트 실패 ⟺ baked_text_score 판정 초과.
        report = aud.audit_native_density(path, tokens=_tokens())
        baked_failures = _slide_baked_failures(report, 1)
        if expected_baked:
            assert baked_failures, (
                f"판정 초과 이미지(kind={kind})는 풀블리드로 채택되면 안 됨"
                f"(베이크 텍스트 실패 보고 기대) — failures={report.failures}")
        else:
            assert not baked_failures, (
                f"판정 미초과 장식(kind={kind})은 풀블리드로 채택돼야 함"
                f"(베이크 텍스트 실패 없어야 함) — {baked_failures}")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        _clean_decor_tmp()


# ===========================================================================
# Property 13 (Task 16.2) — 콘텐츠 텍스트는 겹치는 이미지보다 앞 z-order (Req 6.2, 8.9, 11.4)
# ===========================================================================
# Feature: pptx-native-density-render, Property 13: 장식_배경/장식_비주얼(Vertex 포함)을
# 가진 슬라이드에 대해, 그 이미지와 겹치는 모든 편집가능_네이티브 콘텐츠 텍스트 셰이프는
# 이미지보다 앞선 z-순서(z 가 더 큼)에 배치되어 가려지지 않는다.
@settings(max_examples=100, deadline=None, suppress_health_check=_SUPPRESS)
@given(
    seed=st.integers(min_value=0, max_value=255),
    title=st.sampled_from(_CONTENT_POOL),
    body=st.sampled_from(_CONTENT_POOL),
    img_w=st.floats(min_value=5.0, max_value=8.0),
    img_h=st.floats(min_value=3.5, max_value=5.0),
)
def test_p13_content_text_in_front_of_overlapping_image(seed, title, body, img_w, img_h):
    """emit_decorative_figure 로 장식 이미지를 back-most 로 깐 뒤 콘텐츠 텍스트를 add 하면,
    겹치는 콘텐츠 텍스트는 항상 이미지보다 앞 z-order(spTree 인덱스 더 큼)에 놓인다.
    audit (g) z-order 검사도 위반을 보고하지 않는다.
    """
    fake = _FakeVertexClient(_tiny_png(seed), enabled=True)
    # 이미지(비-풀블리드)와 겹치도록 텍스트를 이미지 영역 안에 배치.
    img_region = (1.0, 1.0, img_w, img_h)
    tmp = tempfile.mkdtemp()
    try:
        with _vertex(fake, optin="1"):
            prs, slide = _new_slide()
            ps = nlr.emit_decorative_figure(
                slide, img_region, "abstract decorative figure, no text",
                region_kind="figure", out_dir=tmp)
            assert ps is not None, "Vertex mock(ON) 시 장식 figure 가 방출돼야 함"
            # 콘텐츠 텍스트 — 이미지와 겹치는 위치에 add(이미지보다 뒤에 add → z 큼).
            nlr.emit_title(slide, title, _tokens(), (1.4, 1.3, img_w - 0.8, 1.0))
            nlr._emit_text_block(slide, body, _tokens(),
                                 (1.4, 2.6, img_w - 0.8, 1.0), role="body")
            path = _save_pptx(prs, tmp)

        pics = _pictures(slide)
        assert len(pics) == 1, "장식 이미지 1개가 임베드돼야 함"
        pic_idx, pic_sh = pics[0]
        pic_rect = _rect_in(pic_sh)
        texts = _content_text_shapes(slide)
        assert len(texts) >= 2, "콘텐츠 텍스트 셰이프가 2개 이상이어야 함"

        overlapped = 0
        for tz, tsh in texts:
            # 겹치는 텍스트는 반드시 이미지보다 앞(z 인덱스 더 큼)이어야 한다.
            if _overlap_ratio(_rect_in(tsh), pic_rect) > 0.0:
                overlapped += 1
                assert tz > pic_idx, (
                    f"겹치는 콘텐츠 텍스트는 이미지보다 앞 z-order 여야 함 "
                    f"(text z={tz}, image z={pic_idx})")
        assert overlapped >= 1, "테스트 전제: 콘텐츠 텍스트가 이미지와 겹쳐야 함"

        # audit (g): z-order 위반 보고 없음.
        report = aud.audit_native_density(path, tokens=_tokens())
        assert not _slide_zorder_failures(report, 1), (
            f"z-order 위반이 없어야 함 — {report.failures}")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        _clean_decor_tmp()


# ===========================================================================
# Property 16 (Task 16.3) — 이미지 손실-0 임베드 + 장식-콘텐츠 경계 불변식 (Req 9.2, 11.4)
# ===========================================================================
# Feature: pptx-native-density-render, Property 16: 슬라이드에 임베드되는 Vertex/명시
# 이미지에 대해, 임베드된 이미지 바이트는 원본과 동일하고(손실 0), 이미지 경계 사각형은
# 슬라이드_경계 안에 있으며, 그 이미지와 각 콘텐츠 텍스트 셰이프 사이의 겹침률은 10%
# 미만이고, 콘텐츠 텍스트 셰이프는 그 이미지보다 앞선 z-순서에 있다.
@settings(max_examples=100, deadline=None, suppress_health_check=_SUPPRESS)
@given(
    seed=st.integers(min_value=0, max_value=255),
    px=st.integers(min_value=1, max_value=64),     # 1x1 PNG 등 다양한 크기
    img_top=st.floats(min_value=0.8, max_value=1.5),
    img_w=st.floats(min_value=4.0, max_value=5.5),
    img_h=st.floats(min_value=3.0, max_value=4.5),
    title=st.sampled_from(_CONTENT_POOL),
)
def test_p16_lossless_embed_and_decoration_content_invariants(
        seed, px, img_top, img_w, img_h, title):
    """Vertex mock 이 생성한 PNG(1x1 포함)를 임베드할 때의 4대 불변식:
    (1) 손실-0(저장 바이트 == 원본 디코드 바이트), (2) 이미지 경계 안,
    (3) 텍스트-이미지 겹침 < 10%, (4) 텍스트가 이미지보다 앞 z-order.
    이미지는 우측, 콘텐츠 텍스트는 좌측으로 분리 배치(겹침<10%).
    """
    raw = _tiny_png(seed, px, px)  # mock 이 반환할 원본 바이트
    fake = _FakeVertexClient(raw, enabled=True)
    img_left = 7.0  # 우측(텍스트 좌측 영역과 분리)
    img_region = (img_left, img_top, img_w, img_h)
    tmp = tempfile.mkdtemp()
    try:
        with _vertex(fake, optin="1"):
            prs, slide = _new_slide()
            nlr.emit_decorative_figure(
                slide, img_region, "abstract decorative figure, no text",
                region_kind="figure", out_dir=tmp)
            # 콘텐츠 텍스트 — 좌측 영역(이미지와 비겹침).
            nlr.emit_title(slide, title, _tokens(), (0.6, 0.8, 5.5, 1.2))
            nlr._emit_text_block(slide, [title, "추가 본문 한 줄"], _tokens(),
                                 (0.6, 2.2, 5.5, 3.0), role="body", bullets=True)

        pics = _pictures(slide)
        assert len(pics) == 1, "이미지 1개가 임베드돼야 함"
        pic_idx, pic_sh = pics[0]

        # (1) 손실-0: 임베드 바이트 == base64 디코드 원본 바이트.
        assert pic_sh.image.blob == raw, "임베드 이미지가 손실-0(원본 바이트 동일)이어야 함"

        # (2) 이미지 경계 안.
        pic_rect = _rect_in(pic_sh)
        assert _within_bounds(pic_rect), f"이미지가 경계 안이어야 함 rect={pic_rect}"

        # (3)/(4) 각 콘텐츠 텍스트와의 겹침 < 10% AND 텍스트가 앞 z-order.
        texts = _content_text_shapes(slide)
        assert len(texts) >= 1, "콘텐츠 텍스트 셰이프가 1개 이상이어야 함"
        for tz, tsh in texts:
            ratio = _overlap_ratio(_rect_in(tsh), pic_rect)
            assert ratio < 0.10, f"텍스트-이미지 겹침률 < 10% 여야 함 (ratio={ratio:.3f})"
            assert tz > pic_idx, (
                f"콘텐츠 텍스트는 이미지보다 앞 z-order 여야 함 "
                f"(text z={tz}, image z={pic_idx})")
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        _clean_decor_tmp()


# ===========================================================================
# Property 17 (Task 16.4) — Vertex 옵트인 OFF 시 콘텐츠 손실 0 (Req 11.3)
# ===========================================================================
# Feature: pptx-native-density-render, Property 17: Vertex 옵트인이 비활성
# (AE_ENABLE_VERTEX_IMAGE != 1 또는 자격증명 부재)인 슬라이드에 대해, 장식_비주얼은
# 생성되지 않지만 입력 콘텐츠 텍스트 런은 전수 보존되어 편집가능_네이티브로 렌더된다
# (콘텐츠 텍스트 손실 0).
@settings(max_examples=100, deadline=None, suppress_health_check=_SUPPRESS)
@given(
    off_mode=st.sampled_from(["env_unset", "env_zero", "client_disabled"]),
    left=st.lists(st.sampled_from(_CONTENT_POOL), min_size=2, max_size=4, unique=True),
    right=st.lists(st.sampled_from(_CONTENT_POOL), min_size=2, max_size=4, unique=True),
)
def test_p17_optin_off_content_lossless_no_decoration(off_mode, left, right):
    """옵트인 OFF 3경로(env 미설정 / env="0" / 자격증명 부재(enabled=False)) 모두:
    maybe_generate_decorative / maybe_add_decorative_background 는 장식을 만들지 않고
    (None / []), Vertex generate 는 호출되지 않으며, render_native_layout 은 콘텐츠
    텍스트 런을 전수 보존하고 이미지를 임베드하지 않는다(콘텐츠 손실 0).
    """
    # off_mode → (env 값, client.enabled)
    if off_mode == "env_unset":
        optin, enabled = None, True
    elif off_mode == "env_zero":
        optin, enabled = "0", True
    else:  # client_disabled — 옵트인 ON 이지만 자격증명 부재(enabled=False)
        optin, enabled = "1", False

    fake = _FakeVertexClient(_tiny_png(7), enabled=enabled)
    data = {"title": _TITLE_TEXT, "subtitle": _SUBTITLE_TEXT,
            "left_content": list(left), "right_content": list(right)}
    tmp = tempfile.mkdtemp()
    try:
        with _vertex(fake, optin=optin):
            # 1) maybe_generate_decorative → None.
            assert nlr.maybe_generate_decorative(
                "decorative prompt", "background", out_dir=tmp) is None, (
                "옵트인 OFF 시 maybe_generate_decorative 는 None 이어야 함")
            # 2) maybe_add_decorative_background → [] (장식 없음).
            _prs0, _slide0 = _new_slide()
            assert nlr.maybe_add_decorative_background(_slide0, data, "cover") == [], (
                "옵트인 OFF 시 maybe_add_decorative_background 는 [] 여야 함")
            # 3) render_native_layout — 콘텐츠 전수 보존 + 장식 없음.
            prs, slide = _new_slide()
            res = nlr.render_native_layout(slide, prs, "two_column", data, _tokens())
            assert res.ok, "two_column 렌더는 성공해야 함"

        # Vertex generate 는 한 번도 호출되지 않아야 한다(네트워크 0·장식 생략).
        assert fake.calls == 0, f"옵트인 OFF 시 Vertex generate 호출 0 이어야 함 (calls={fake.calls})"
        # 장식 이미지 없음.
        assert len(_pictures(slide)) == 0, "옵트인 OFF 시 장식 이미지가 없어야 함"
        # 콘텐츠 텍스트 런 전수 보존(손실 0).
        all_text = _all_editable_text(slide)
        for s in [_TITLE_TEXT, _SUBTITLE_TEXT, *left, *right]:
            assert s in all_text, f"콘텐츠 텍스트 '{s}' 가 편집가능 네이티브로 보존돼야 함"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        _clean_decor_tmp()


# ===========================================================================
# Property 18 (Task 16.5) — Vertex 장식 슬라이드의 베이크 미초과 + 콘텐츠 공존 (Req 11.6)
# ===========================================================================
# Feature: pptx-native-density-render, Property 18: Vertex 로 생성된 장식_비주얼을
# 포함하는 슬라이드에 대해, audit_native_density 는 그 장식_비주얼이 baked_text_score
# 판정(비율 < 6% AND 줄 < 6)을 초과하지 않음을 확인하고, 동일 슬라이드에 비어있지 않은
# 텍스트 런을 가진 편집가능_네이티브 콘텐츠 셰이프가 1개 이상 공존함을 확인한다.
@settings(max_examples=100, deadline=None, suppress_health_check=_SUPPRESS)
@given(
    layout=st.sampled_from(["cover", "section_divider"]),
    seed=st.integers(min_value=0, max_value=255),
)
def test_p18_vertex_decoration_not_baked_and_content_coexists(layout, seed):
    """Vertex mock(추상 장식_배경, 텍스트 없음)을 풀블리드로 포함한 슬라이드에서,
    audit_native_density 는 (f) 베이크 텍스트 판정 미초과를 확인하고, 동일 슬라이드에
    편집가능 콘텐츠 텍스트 셰이프가 1개 이상 공존한다(Req 11.6).
    """
    png = _abstract_png(seed)
    assert not _baked_exceeded(png), "전제: 추상 장식은 baked 판정을 미초과해야 함"

    fake = _FakeVertexClient(png, enabled=True)
    data = {"title": _TITLE_TEXT, "eyebrow": "EYEBROW", "subtitle": _SUBTITLE_TEXT,
            "footer": "2026", "section_number": 1,
            "description": "섹션 설명 문장"}
    tmp = tempfile.mkdtemp()
    try:
        with _vertex(fake, optin="1"):
            prs, slide = _new_slide()
            res = nlr.render_native_layout(slide, prs, layout, data, _tokens())
            assert res.ok, f"{layout} 렌더는 성공해야 함"
            path = _save_pptx(prs, tmp)

        # 장식_배경(Vertex)이 풀블리드로 임베드되었는지 확인.
        pics = _pictures(slide)
        assert len(pics) == 1, "Vertex 장식_배경 1개가 임베드돼야 함"

        report = aud.audit_native_density(path, tokens=_tokens())
        # (f) 장식_비주얼이 베이크 텍스트 판정 미초과.
        assert not _slide_baked_failures(report, 1), (
            f"Vertex 장식은 베이크 텍스트 판정을 초과하지 않아야 함 — {report.failures}")
        # 콘텐츠 공존: 편집가능 콘텐츠 텍스트 셰이프 ≥ 1.
        content = _content_text_shapes(slide)
        assert len(content) >= 1, "편집가능 콘텐츠 텍스트 셰이프가 1개 이상 공존해야 함"
        assert _TITLE_TEXT in _all_editable_text(slide), "제목 콘텐츠가 공존해야 함"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
        _clean_decor_tmp()


if __name__ == "__main__":  # pragma: no cover
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
