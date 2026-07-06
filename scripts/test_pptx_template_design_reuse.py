"""Regression — 템플릿의 디자인(배경·이미지·장식)을 *복사*해 새 내용을 채운다(Genspark 방식).

사용자 보고: "템플릿 슬라이드에 새로운 슬라이드만 추가하여 내용이 작성됨. 템플릿은 디자인
양식·배경·스타일을 카피하여 거기에 새 내용을 작성하는 것."

수정: `_tool_generate_pptx`가 템플릿이면 템플릿의 디자인 슬라이드를 "도너"로 재사용하고,
도너보다 많이 필요하면 `_clone_slide`로 디자인(배경+이미지+도형)을 복제한 뒤 텍스트만
새 내용으로 교체한다. 잉여 도너는 제거한다.

Correctness properties:
  P1. 템플릿 디자인 슬라이드의 배경/이미지가 모든 출력 슬라이드에 보존된다(디자인 복사).
  P2. 출력 슬라이드 수 == 1(cover) + len(slides) (도너 잉여분 제거, 샘플 미혼입).
  P3. 템플릿의 원본 샘플 텍스트는 결과물에 남지 않고, 새 제목/내용으로 교체된다.
  P4. _clone_slide는 배경과 이미지(rId)를 보존한다.

실행: pytest scripts/test_pptx_template_design_reuse.py -q
"""
from __future__ import annotations

import os
import sys
import io
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


_PNG_1x1 = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082"
)


def _make_designed_template(path, n_slides=1, marker_prefix="SAMPLE"):
    """배경 채우기 + 이미지 + 제목/본문 placeholder를 가진 '디자인 템플릿' pptx 생성."""
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    for i in range(n_slides):
        s = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
        s.shapes.title.text = f"{marker_prefix}_TITLE_{i}"
        if len(s.placeholders) > 1:
            s.placeholders[1].text = f"{marker_prefix}_BODY_{i}"
        # 배경 채우기
        cSld = s._element.find(qn('p:cSld'))
        bg = cSld.makeelement(qn('p:bg'), {})
        bgPr = bg.makeelement(qn('p:bgPr'), {})
        sf = bgPr.makeelement(qn('a:solidFill'), {})
        sf.append(sf.makeelement(qn('a:srgbClr'), {'val': '0B5394'}))
        bgPr.append(sf); bgPr.append(bgPr.makeelement(qn('a:effectLst'), {}))
        bg.append(bgPr); cSld.insert(0, bg)
        # 장식 이미지(로고 대용)
        s.shapes.add_picture(io.BytesIO(_PNG_1x1), Inches(11.5), Inches(0.3), Inches(1), Inches(1))
    prs.save(str(path))


def _analyze(pptx_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    prs = Presentation(pptx_path)
    out = {"count": len(prs.slides._sldIdLst), "slides": [], "all_text": ""}
    texts = []
    for sl in prs.slides:
        pics = sum(1 for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
        has_bg = sl._element.find(qn('p:cSld')).find(qn('p:bg')) is not None
        out["slides"].append({"pictures": pics, "has_bg": has_bg})
        for sh in sl.shapes:
            if sh.has_text_frame:
                texts.append(sh.text_frame.text)
    out["all_text"] = "\n".join(texts)
    return out


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_clone_slide_preserves_bg_and_image(tmp_path):
    """P4 — _clone_slide가 배경+이미지를 보존."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    tpl = tmp_path / "t.pptx"
    _make_designed_template(tpl, n_slides=1)
    prs = Presentation(str(tpl))
    src = list(prs.slides)[0]
    clone = server._clone_slide(prs, src)
    assert clone is not None
    pics = sum(1 for sh in clone.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
    assert pics == 1, "복제본에 이미지가 보존되지 않음"
    assert clone._element.find(qn('p:cSld')).find(qn('p:bg')) is not None, "복제본 배경 손실"


def test_template_design_copied_to_all_slides(tmp_path):
    """P1+P2+P3 — 단일 디자인 도너로 3개 콘텐츠 생성: 디자인 복사 + 샘플 텍스트 교체."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    tpl = tmp_path / "designed.pptx"
    _make_designed_template(tpl, n_slides=1, marker_prefix="ORIGSAMPLE")

    slides = [
        {"title": "소개", "bullets": ["새 내용 A", "새 내용 B"]},
        {"title": "현황", "bullets": ["새 내용 C"]},
        {"title": "결론", "bullets": ["새 내용 D"]},
    ]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "새 발표", "slides": slides, "templatePath": str(tpl)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        st = _analyze(res["absPath"])
        # P2 — cover(1) + content(3) = 4 (도너 1개 재사용 + 3개 복제, 잉여 없음)
        assert st["count"] == 4, f"슬라이드 수 {st['count']} (기대 4)"
        # P1 — 모든 슬라이드가 템플릿 배경 + 디자인 이미지를 보존
        for i, sl in enumerate(st["slides"]):
            assert sl["has_bg"], f"슬라이드 {i} 배경 손실(디자인 미복사)"
            assert sl["pictures"] >= 1, f"슬라이드 {i} 디자인 이미지 손실"
        # P3 — 원본 샘플 텍스트 제거 + 새 제목/내용 존재
        assert "ORIGSAMPLE" not in st["all_text"], "템플릿 샘플 텍스트가 남음"
        assert "새 발표" in st["all_text"], "새 표지 제목 없음"
        assert "소개" in st["all_text"] and "결론" in st["all_text"], "새 슬라이드 제목 누락"
    finally:
        _cleanup(res)


def test_template_excess_donors_removed(tmp_path):
    """P2 — 콘텐츠 수 < 도너 수: 잉여 도너 슬라이드 제거(샘플 미혼입)."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    tpl = tmp_path / "many.pptx"
    _make_designed_template(tpl, n_slides=5, marker_prefix="EXTRA")  # 도너 5개

    slides = [{"title": "한 장", "bullets": ["x"]}]  # cover + 1 content = 2
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "짧은 덱", "slides": slides, "templatePath": str(tpl)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        st = _analyze(res["absPath"])
        assert st["count"] == 2, f"슬라이드 수 {st['count']} (기대 2 — 잉여 도너 3개 제거)"
        assert "EXTRA" not in st["all_text"], "잉여 도너 샘플 텍스트가 남음"
        assert "짧은 덱" in st["all_text"]
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
