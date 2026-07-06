"""슬라이드별 모든 PICTURE를 z-order 순서로 분류 — 무엇이 올라가 있는지 규명.

각 이미지를 분류: solid(단색)/gradient/photo-illustration/text-baked/thin-bar/icon.
풀블리드 여부, 위치, 크기, 종횡비, 고유색, 엣지밀도로 판정.
"""
from __future__ import annotations
import io, sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

EMU = 914400.0
SW, SH = 13.333, 7.5

def _in(v): return round(float(v)/EMU, 2) if v is not None else None
def _rect(sh):
    try:
        l,t,w,h = sh.left, sh.top, sh.width, sh.height
        if None in (l,t,w,h): return None
        return (_in(l),_in(t),_in(w),_in(h))
    except Exception: return None
def _fb(r): return r and r[0]<=0.3 and r[1]<=0.3 and r[2]>=SW*0.92 and r[3]>=SH*0.92

def classify(blob):
    try:
        im = Image.open(io.BytesIO(blob)).convert("RGB")
    except Exception as e:
        return f"디코드실패({e.__class__.__name__})", {}
    w,h = im.size
    small = im.resize((min(w,360), min(h,202)))
    px = list(small.getdata())
    uniq = len(set(px))
    # 엣지밀도
    g = [(r*299+gg*587+b*114)//1000 for r,gg,b in px]
    sw = small.size[0]; edges=tot=0
    for y in range(small.size[1]):
        base=y*sw
        for x in range(sw-1):
            tot+=1
            if abs(g[base+x]-g[base+x+1])>40: edges+=1
    epct = edges/tot*100 if tot else 0
    ar = w/h if h else 0
    meta = {"size": f"{w}x{h}", "uniq": uniq, "edge": round(epct,1), "ar": round(ar,1)}
    if w<8 or h<8: return f"초소형 {w}x{h}", meta
    if ar>15 or ar<0.07: return f"thin-bar(종횡비{ar:.0f}) {w}x{h}", meta
    if uniq<=4: return "solid(단색)", meta
    if uniq<60 and epct<2: return "gradient/단순", meta
    if epct>=8 and uniq>=2000: return "text-baked/도표(텍스트구워짐 의심)", meta
    if uniq>=1500: return "photo/illustration", meta
    return "기타", meta

def main(path):
    prs = Presentation(path)
    print(f"=== 미디어 분류(z-order별): {path}\n")
    for idx, slide in enumerate(prs.slides,1):
        print(f"[슬라이드 {idx}]")
        for z, sh in enumerate(slide.shapes):
            try: is_pic = sh.shape_type==MSO_SHAPE_TYPE.PICTURE
            except Exception: is_pic=False
            if not is_pic: continue
            r=_rect(sh)
            try: blob=sh.image.blob
            except Exception: blob=None
            kind,meta = classify(blob) if blob else ("blob없음",{})
            fb = "풀블리드" if _fb(r) else "부분"
            print(f"   z={z} {fb} @{r} → {kind}  {meta}")
        print()

if __name__=="__main__":
    main(sys.argv[1] if len(sys.argv)>1 else "")
