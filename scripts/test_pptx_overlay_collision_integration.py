"""Hermetic integration test — spec: pptx-overlay-collision-fix (bugfix), Task 6.

Drives the REAL native placement path (``ai_engine.native_diagram_pptx.
build_native_cover`` / ``build_native_diagram``) end to end over a MIXED deck and
asserts — at the *deck* level — that the overlay/collision defects A / B / C are
resolved (design "Integration Tests" / Property 1·2):

  (a) DEFECT A·C resolved — for every slide in the deck, every
      (textBox ∪ badge) pair overlaps less than 10% of the smaller box
      (``overlap_area(a,b) < 0.10*min(area(a),area(b))``). This covers the
      cover title↔subtitle vertical stack (defect A) and the number badge↔label
      gutter placement (defect C) simultaneously, measured with the EXISTING
      audit tool ``scripts/audit_pptx_textbox_overlap.py`` (``boxes`` / ``ov``).

  (b) DEFECT B + loss-zero — the pure decision seam
      ``ai_engine.server._select_render_plan(role="structural",
      has_vertex_image=True, bg_has_baked_text=True)`` returns
      ``primary == "NATIVE_SHAPES"`` (structural stays editable native shapes,
      not a baked-text full-bleed body carrier) ∧ ``vertex_slot == "backdrop"``
      (the generated Vertex image is preserved — loss-zero) ∧
      ``body_separated is True`` (body/background separation signalled).

The mixed deck assembles four slides through the real code:
  (a) cover    — long title + subtitle  → ``build_native_cover`` (defect A path)
  (b) twocol   — 6 numbered-badge items → ``build_native_diagram("twocol", …)`` (defect C path)
  (c) structural flow                   → ``build_native_diagram("flow", …)``   (defect B path)
  (d) block diagram                     → ``build_native_diagram("block", …)``  (defect C path)

Everything is hermetic — no network. Only the native assembly functions and the
pure decision function ``_select_render_plan`` are exercised. The Bedrock
gateway, the Vertex client and the HTML→PNG renderer are never invoked. Overlap
is measured with the EXISTING audit tool's axis-aligned intersection
(``audit_pptx_textbox_overlap.ov`` / ``boxes``) — the same definition used by the
fix code's ``layout_geometry.overlap_area`` (audit ↔ code parity). Helper
patterns (``_blank_169_presentation`` etc.) mirror
``test_pptx_overlay_collision_bug_condition.py``.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_overlay_collision_integration.py -p no:cacheprovider -q

_Requirements: 2.1, 2.2, 2.3, 2.4, 2.5_
"""
from __future__ import annotations

import os
import sys

# Make ai_engine (repo root) and the audit tools (scripts/) importable.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# Native assembly path (real code under test) — no network.
from ai_engine.native_diagram_pptx import (  # noqa: E402
    build_native_cover,
    build_native_diagram,
)
# Pure decision function (real code under test) — no network.
from ai_engine.server import _select_render_plan  # noqa: E402

# EXISTING audit measurement functions (reused — same axis-aligned intersection).
import audit_pptx_textbox_overlap as tov  # noqa: E402


# Threshold of "meaningful overlap" — design Bug Condition (10% of smaller box).
THRESHOLD = 0.10
# 16:9 slide so the documented inch coordinates hold (design Examples).
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ──────────────────────────────────────────────────────────────────────────
# Helpers — mirror test_pptx_overlay_collision_bug_condition.py
# ──────────────────────────────────────────────────────────────────────────
def _blank_169_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs):
    # layout[6] = fully blank layout (no placeholders to interfere).
    return prs.slides.add_slide(prs.slide_layouts[6])


def _area(box):
    # box = (head, l, t, w, h) — same shape audit tools use.
    return max(0.0, box[3]) * max(0.0, box[4])


def _is_badge(box):
    """A number badge: very short purely-numeric label + small box."""
    head = (box[0] or "").strip()
    return head.isdigit() and 1 <= len(head) <= 2


def _worst_pair_overlap(boxes):
    """Return (ratio, info) for the pair with the largest overlap / min(area)."""
    worst = 0.0
    info = None
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ov = tov.ov(boxes[i], boxes[j])
            if ov <= 0.0:
                continue
            amin = min(_area(boxes[i]), _area(boxes[j]))
            if amin <= 0.0:
                continue
            ratio = ov / amin
            if ratio > worst:
                worst = ratio
                info = (boxes[i], boxes[j], ov, amin)
    return worst, info


def _assert_slide_no_overlap(slide, label):
    """Assert every (textBox ∪ badge) pair on a slide overlaps < 10% min(area)."""
    boxes = tov.boxes(slide)
    assert len(boxes) >= 2, f"[{label}] 텍스트 박스가 2개 미만(전제 실패): {boxes}"
    worst, info = _worst_pair_overlap(boxes)
    detail = ""
    if info is not None:
        a, b, ov, amin = info
        detail = (f"  '{a[0]}'@({a[1]},{a[2]} {a[3]}x{a[4]}) ↔ "
                  f"'{b[0]}'@({b[1]},{b[2]} {b[3]}x{b[4]})  "
                  f"overlap={ov:.2f}in² ({worst*100:.0f}% of smaller {amin:.2f}in²)")
    assert worst < THRESHOLD, (
        f"[{label}] 결함 A/C 재발 — (텍스트박스 ∪ 배지) 쌍 겹침 "
        f"{worst*100:.0f}% ≥ 임계 {THRESHOLD*100:.0f}%.\n" + detail
    )
    return boxes


# ──────────────────────────────────────────────────────────────────────────
# Mixed deck — assembled with the REAL native placement path (no network).
# ──────────────────────────────────────────────────────────────────────────
def _assemble_mixed_deck():
    """Assemble cover + twocol + structural-flow + block into one in-memory deck.

    Returns the Presentation plus the four slides for per-slide inspection."""
    prs = _blank_169_presentation()

    # (a) 표지 — 긴 제목 + 부제 (결함 A 경로: 제목 박스 하단 > 부제 top 고정 오프셋).
    cover = _blank_slide(prs)
    drew_cover = build_native_cover(
        cover, prs,
        title="데이터 흐름 구조 정리",
        subtitle="조직 전반의 디렉토리 깊이와 데이터 흐름을 한눈에 시각화합니다",
        eyebrow="", date_str="", kpis=None,
    )
    assert drew_cover, "build_native_cover 가 표지를 그리지 못함(전제 실패)"

    # (b) twocol — 번호 배지 6항목 (결함 C 경로: 배지가 라벨 카드 내부에 포개짐).
    twocol = _blank_slide(prs)
    twocol_content = (
        "언어: Python 3.11\n"
        "프레임워크: FastAPI\n"
        "런타임: Electron\n"
        "빌드: webpack\n"
        "테스트: pytest\n"
        "CI: GitHub Actions"
    )
    drew_twocol = build_native_diagram(
        twocol, "twocol", twocol_content,
        region=(0.6, 1.7, 12.1, 5.2), palette=None, title="기술 스택",
    )
    assert drew_twocol, "build_native_diagram(twocol) 가 그리지 못함(전제 실패)"

    # (c) structural flow — 흐름도 (결함 B 경로: 구조형은 편집 가능 네이티브 도형 유지).
    flow = _blank_slide(prs)
    flow_content = (
        "데이터 수집: 외부 소스에서 원천 데이터를 적재\n"
        "정제: 결측/중복 제거 및 표준화\n"
        "분석: 지표 산출 및 패턴 탐지\n"
        "리포팅: 대시보드로 결과 시각화"
    )
    drew_flow = build_native_diagram(
        flow, "flow", flow_content,
        region=(0.6, 1.7, 12.1, 5.2), palette=None, title="데이터 파이프라인",
    )
    assert drew_flow, "build_native_diagram(flow) 가 그리지 못함(전제 실패)"

    # (d) block — 번호 배지 리스트 (결함 C 경로: 좌측 거터 배지).
    block = _blank_slide(prs)
    block_content = (
        "보안: 전 구간 암호화와 접근 통제\n"
        "확장성: 수평 확장 가능한 무상태 서비스\n"
        "관측성: 로그·메트릭·트레이스 통합\n"
        "복원력: 장애 격리와 자동 복구"
    )
    drew_block = build_native_diagram(
        block, "block", block_content,
        region=(0.6, 1.7, 12.1, 5.2), palette=None, title="설계 원칙",
    )
    assert drew_block, "build_native_diagram(block) 가 그리지 못함(전제 실패)"

    return prs, cover, twocol, flow, block


# ==========================================================================
# (a)(C) DEFECT A·C resolved at the deck level — no meaningful text/badge overlap
# ==========================================================================
def test_integration_mixed_deck_no_textbox_badge_overlap():
    """혼합 덱(표지·twocol·flow·block)을 실제 네이티브 경로로 조립한 뒤, 각 슬라이드의
    모든 (텍스트박스 ∪ 배지) 쌍이 겹침 < 10% min(area) 임을 audit 도구로 단언한다.
    표지 제목↔부제(결함 A)와 번호 배지↔라벨(결함 C)이 덱 레벨에서 해소됐음을 확인한다."""
    prs, cover, twocol, flow, block = _assemble_mixed_deck()

    # 슬라이드별 (텍스트박스 ∪ 배지) 비겹침 — 결함 A·C 덱 레벨 검증.
    _assert_slide_no_overlap(cover, "표지(cover)")
    twocol_boxes = _assert_slide_no_overlap(twocol, "twocol")
    _assert_slide_no_overlap(flow, "structural-flow")
    block_boxes = _assert_slide_no_overlap(block, "block")

    # 배지가 실제로 존재함을 확인(결함 C 경로가 구동됐다는 전제) — twocol/block.
    twocol_badges = [b for b in twocol_boxes if _is_badge(b)]
    block_badges = [b for b in block_boxes if _is_badge(b)]
    assert twocol_badges, f"twocol 번호 배지를 찾지 못함(전제 실패): {[b[0] for b in twocol_boxes]}"
    assert block_badges, f"block 번호 배지를 찾지 못함(전제 실패): {[b[0] for b in block_boxes]}"

    # 표지 텍스트 박스(제목+부제) — 결함 A 대상 쌍이 실제로 존재함을 확인.
    cover_boxes = tov.boxes(cover)
    assert len(cover_boxes) >= 2, f"표지 텍스트 박스가 2개 미만: {cover_boxes}"


# ==========================================================================
# (b) DEFECT B + loss-zero — structural baked-text decision at the deck level
# ==========================================================================
def test_integration_structural_baked_text_native_carrier():
    """덱 레벨 보조 단언 — 텍스트가 구워진 풀블리드 배경이 깔린 구조형 슬라이드에서
    ``_select_render_plan`` 이 그 이미지를 본문 캐리어로 쓰지 않고(NATIVE_SHAPES),
    생성된 Vertex 이미지는 backdrop(손실-0)으로 보존하며, 본문/배경 분리를 신호함을
    단언한다(결함 B + 손실-0)."""
    plan = _select_render_plan(
        role="structural",
        has_vertex_image=True,
        has_native_diagram=False,
        has_image_file=False,
        has_slide_bg=False,
        html_enabled=False,
        bg_has_baked_text=True,
    )
    assert plan["primary"] == "NATIVE_SHAPES", (
        "결함 B — 구조형이 편집 네이티브 도형 대신 "
        f"{plan['primary']} 로 렌더된다(구워진-텍스트 이미지가 본문 캐리어)."
    )
    assert plan["vertex_slot"] == "backdrop", (
        "결함 B — 생성 Vertex 이미지가 backdrop(손실-0 보존)이 아니라 "
        f"{plan['vertex_slot']!r} 슬롯에 놓인다."
    )
    assert plan["body_separated"] is True, (
        "결함 B — 본문/배경 분리 의도(body_separated)가 신호되지 않는다."
    )


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
