"""Bug-condition exploration test — spec: pptx-fullbleed-native-overlay-collision
(bugfix), Task 1.

PROPERTY 1 — Bug Condition (풀블리드-네이티브 오버레이 충돌 + 경계밖 도형 재현).

This is an EXPLORATORY bug-condition test. It drives the REAL
``ai_engine.server._tool_generate_pptx`` over the audited "템플릿 + 이미지 슬라이드"
path — the ONLY path where the defect lives (per bugfix.md scope note): a slide
gets BOTH a content-baked full-bleed HTML background (``slideBackground``) AND the
same content as native title/body shapes. The test encodes the EXPECTED (post-fix)
behaviour (design "expectedBehavior"):

  (1) every native text shape overlaps the content-baked full-bleed background by
      < 10% of its own area,
  (2) the slide title is represented exactly ONCE (not baked-in AND native),
  (3) ALL shapes lie inside the slide bounds (0,0,13.333,7.5) — no negative-top /
      off-slide shapes.

On the UNFIXED code these are violated — the native body TEXT_BOX (0.6,1.6,12.1,5.4)
sits on top of the baked full-bleed (~100% overlap), the title appears both baked
into the PNG and as a native TEXT_BOX (0.6,0.3,12.1,1.0), and the native cover adds
a decorative OVAL at top=-1.5 (off-slide). Those FAILURES are the proof the bug
exists. After the fix (tasks 3.x) the SAME test is re-run (task 3.4) and must PASS.

DO NOT "fix" this test or the code when it fails on unfixed code — the failure is
intended; it surfaces the counterexamples that prove the defect.

How the audited dual layer is reproduced hermetically (no network):
  * A 16:9 TEMPLATE deck (templatePath) → ``used_template=True`` so the pipeline
    reuses donor slides and STRIPS them to background. With no title/body
    placeholder, ``_safe_set_title`` falls back to add_textbox(0.6,0.3,12.1,1.0)
    and the body fallback adds add_textbox(0.6,1.6,12.1,5.4) — the exact audited
    native shapes (server.py _safe_set_title 4018 / body fallback 5174).
  * HTML slides ON + ``_generate_html_slide_for_section`` MOCKED to BAKE the
    slide title + body into the returned full-bleed PNG → the pipeline sets
    ``slideBackground`` to that baked PNG (server.py HTML gate 5126/5149).
  * Vertex hero/body backgrounds DISABLED (env) → no coverBackground → the native
    cover is built with ``over_image=False`` → decorative OVAL at top=-1.5
    (native_diagram_pptx.py 1597), reproducing the off-slide shape.

Measurement reuses the EXISTING audit tools:
  * ``scripts/audit_pptx_textbox_overlap.py`` (``ov`` / ``_in``) — axis-aligned
    area overlap (native text ↔ baked full-bleed background),
  * ``scripts/audit_pptx_baked_text.py`` (``baked_text_score``) — confirms the
    full-bleed background really carries baked text (so the title is baked in),
  * ``scripts/audit_pptx_zorder_break.py`` (``_rect`` / ``SW`` / ``SH``) — off-slide
    (out-of-bounds) detection.

Everything is hermetic — the Bedrock gateway (``_get_gw`` / ``_call_bridge``), the
Vertex client (``get_vertex_image_client``), ``_tool_generate_image`` and the
HTML→PNG renderers (``_render_html_slide_to_png`` /
``_generate_html_slide_for_section``) are all mocked. Chrome is never used
(``_find_local_chrome`` → "").

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_fullbleed_native_overlay_bug_condition.py -p no:cacheprovider -q

EXPECTED OUTCOME on UNFIXED code: FAIL (overlap ~100%, title twice, off-slide
OVAL top=-1.5 present).

_Bug_Condition: isBugCondition(slide) = collision OR offSlide (design Bug Condition)_
_Requirements: 1.1, 1.2, 1.3, 1.4_
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile
import zipfile
import io
from unittest.mock import patch

# Make ai_engine (repo root) and the audit tools (scripts/) importable.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402

# EXISTING audit measurement functions (reused — same axis-aligned geometry).
import audit_pptx_textbox_overlap as tov  # noqa: E402
import audit_pptx_baked_text as abt  # noqa: E402
import audit_pptx_zorder_break as azb  # noqa: E402


# Design Bug Condition / expectedBehavior threshold (10% of the native box).
THRESHOLD = 0.10
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Curated content deck (deterministic). Each slide carries a title + bullets so
# the pipeline emits BOTH a baked full-bleed background AND native title/body.
_DECK_TITLE = "프로젝트 구조 분석"
_CONTENT_SLIDES = [
    {"title": "Depth 0 — 루트 구조",
     "bullets": ["cgjang 루트 디렉토리", "주요 모듈 배치", "빌드 산출물 경로"]},
    {"title": "Depth 1 — 코어 모듈",
     "bullets": ["ai_engine 서버", "electron 브리지", "src 프론트엔드"]},
    {"title": "Depth 2 — 데이터 흐름",
     "bullets": ["입력 수집", "정제 및 적재", "분석 파이프라인", "결과 시각화"]},
]


# --------------------------------------------------------------------------
# Hermetic fakes
# --------------------------------------------------------------------------
def _make_baked_section_png(path: str, heading: str, bullets, w=1920, h=1080) -> str:
    """Render a full-bleed section PNG with the slide title + body BAKED into the
    raster (simulates the Genspark-class HTML→PNG render whose Korean title/body
    are pixels, not editable text). High-contrast text rows trigger the baked-text
    detector (``audit_pptx_baked_text.baked_text_score``)."""
    im = Image.new("RGB", (w, h), (24, 38, 70))
    d = ImageDraw.Draw(im)
    # Title baked near the top (mirrors the native title TEXT_BOX content).
    d.text((60, 48), heading, fill=(245, 248, 255))
    d.text((61, 49), heading, fill=(245, 248, 255))  # faux-bold for stronger edges
    # Body bullets baked below — many high-contrast rows across the full width.
    base = ("• " + "   ".join(str(b) for b in bullets)
            + "   0123456789 ABCDEFGH 데이터 정제 적재 분석 파이프라인 ROOT SRC DIST")
    for k in range(16):
        y = 180 + k * 52
        d.text((60, y), base, fill=(232, 238, 250))
    im.save(path)
    return path


def _make_compositing_section_fake(project_path: str, baked_titles: dict):
    """Stand-in for ``_generate_html_slide_for_section`` — bakes heading+body into
    the returned full-bleed PNG and records the baked title per call so the test
    can prove the title is baked into the background (title-twice defect)."""
    state = {"n": 0}

    async def _section_fake(gw, model_id, heading, body, doc_context,
                            proj_path, style_profile=None,
                            hero_image="", render_info=None):
        state["n"] += 1
        gen = os.path.join(project_path, ".generated")
        os.makedirs(gen, exist_ok=True)
        name = f"baked-section-{state['n']}.png"
        out = os.path.join(gen, name)
        # Derive bullets from the body text passed in (one per line).
        bullets = [ln.strip(" •-\t") for ln in str(body).splitlines() if ln.strip()]
        _make_baked_section_png(out, str(heading), bullets or [str(body)[:40]])
        baked_titles[str(heading).strip()] = f".generated/{name}"
        if isinstance(render_info, dict):
            render_info["layout"] = "feature_grid"
            render_info["composited"] = False
        return f".generated/{name}"

    return _section_fake


async def _render_html_png_fake(html, output_path, width=1920, height=1080, timeout=30, **_k):
    """Stand-in for ``_render_html_slide_to_png`` (cover) — writes a valid PNG."""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    Image.new("RGB", (320, 180), (30, 30, 30)).save(output_path)
    return {"ok": True}


async def _img_gen_disabled(*_a, **_k):
    """Stand-in for ``_tool_generate_image`` — never returns a path (no network)."""
    return json.dumps({"error": "test-disabled"})


class _FakeVertexClient:
    """Defensive Vertex stub — Vertex is also disabled via env, so it must never
    actually be called; this guarantees zero network even if a path slips."""

    def __init__(self) -> None:
        self.enabled = False
        self.calls = 0

    async def generate(self, *_a, **_k):
        self.calls += 1
        raise AssertionError("Vertex must not be called in this hermetic test")


# --------------------------------------------------------------------------
# Template builder — a 16:9 deck with donor slides (used_template path).
# --------------------------------------------------------------------------
def _build_template(path: str, n_slides: int = 6) -> str:
    """Create a minimal 16:9 .pptx with N donor slides so the pipeline takes the
    template path (``used_template=True``) and strips donors to background — the
    audited scenario that yields the native title/body TEXT_BOX fallbacks."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    layout = prs.slide_layouts[1]  # Title and Content (placeholders get stripped)
    for k in range(n_slides):
        s = prs.slides.add_slide(layout)
        try:
            s.shapes.title.text = f"도너 {k + 1}"
        except Exception:
            pass
    prs.save(path)
    return path


# --------------------------------------------------------------------------
# Audit helpers (reuse existing tools' geometry)
# --------------------------------------------------------------------------
def _fullbleed_picture_boxes(slide):
    """Return (name,l,t,w,h) boxes for full-bleed PICTURE shapes (the baked
    background), using the audit's inch conversion (``tov._in``)."""
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        l, t, w, h = tov._in(sh.left), tov._in(sh.top), tov._in(sh.width), tov._in(sh.height)
        if None in (l, t, w, h):
            continue
        if l <= 0.3 and t <= 0.3 and w >= SLIDE_W_IN * 0.92 and h >= SLIDE_H_IN * 0.92:
            out.append(("FULLBLEED", l, t, w, h))
    return out


def _picture_is_baked(slide, pic_box) -> bool:
    """Confirm a full-bleed picture carries baked text via baked_text_score."""
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            blob = sh.image.blob
        except Exception:
            continue
        try:
            pct, lines = abt.baked_text_score(Image.open(io.BytesIO(blob)))
        except Exception:
            continue
        if pct >= 6.0 or lines >= 6:
            return True
    return False


def _native_text_boxes(slide):
    """Native (non-picture) text boxes via the existing audit extractor."""
    return tov.boxes(slide)


def _area(box):
    return max(0.0, box[3]) * max(0.0, box[4])


def _worst_text_over_fullbleed(slide):
    """Max overlap ratio (overlap / native-box area) of any native text box over
    any baked full-bleed background. Returns (ratio, detail)."""
    fbs = _fullbleed_picture_boxes(slide)
    if not fbs:
        return 0.0, None
    worst = 0.0
    detail = None
    for tb in _native_text_boxes(slide):
        ta = _area(tb)
        if ta <= 0.0:
            continue
        for fb in fbs:
            ov = tov.ov(tb, fb)
            if ov <= 0.0:
                continue
            ratio = ov / ta
            if ratio > worst:
                worst = ratio
                detail = (f"'{tb[0]}'@({tb[1]},{tb[2]} {tb[3]}x{tb[4]}) over baked "
                          f"full-bleed → {ov:.2f}in² = {ratio * 100:.0f}% of native box")
    return worst, detail


def _title_occurrences(slide, title: str):
    """Count how many times the slide title is represented:
       (# native text shapes whose first line == title) + (1 if a baked full-bleed
       background — which carries the title — is present)."""
    title_n = title.strip()
    native = 0
    for tb in _native_text_boxes(slide):
        if (tb[0] or "").strip() == title_n[:24].strip():
            native += 1
    fbs = _fullbleed_picture_boxes(slide)
    baked = 1 if (fbs and _picture_is_baked(slide, fbs[0])) else 0
    return native + baked, native, baked


def _offslide_shapes(slide):
    """All shapes whose rect falls outside slide bounds, reusing the audit's
    geometry (``audit_pptx_zorder_break._rect`` + its off-slide predicate / SW,SH)."""
    out = []
    for sh in slide.shapes:
        r = azb._rect(sh)
        if not r:
            continue
        if (r[0] < -0.05 or r[1] < -0.05
                or r[0] + r[2] > azb.SW + 0.05 or r[1] + r[3] > azb.SH + 0.05):
            try:
                nm = sh.shape_type
            except Exception:
                nm = "?"
            out.append((str(nm), r))
    return out


# --------------------------------------------------------------------------
# Drive the REAL pipeline once; cache the produced deck for all assertions.
# --------------------------------------------------------------------------
_GENERATED = {"pptx": None, "baked_titles": None}


def _generate_deck():
    if _GENERATED["pptx"] and os.path.isfile(_GENERATED["pptx"]):
        return _GENERATED["pptx"], _GENERATED["baked_titles"]

    proj = tempfile.mkdtemp(prefix="fbno_")
    tpl = _build_template(os.path.join(proj, "template_169.pptx"), n_slides=6)
    baked_titles: dict = {}
    fake_vx = _FakeVertexClient()

    env = {
        "AE_ENABLE_HTML_SLIDES": "1",        # HTML primary → baked full-bleed bg
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_DISABLE_VERTEX_HERO": "1",       # no cover bg → native cover off-slide oval
        "AE_DISABLE_VERTEX_BODY_BG": "1",    # fully hermetic (no Vertex pre-gen)
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {
        "title": _DECK_TITLE,
        "templatePath": tpl,
        "slides": [dict(s) for s in _CONTENT_SLIDES],
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake_vx), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", lambda *a, **k: object()), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_render_html_slide_to_png", _render_html_png_fake), \
            patch.object(srv, "_generate_html_slide_for_section",
                         _make_compositing_section_fake(proj, baked_titles)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))

    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패(전제): {result}"
    _GENERATED["pptx"] = result["absPath"]
    _GENERATED["baked_titles"] = baked_titles
    return result["absPath"], baked_titles


def _content_slides_with_baked_fullbleed(prs):
    """Indices (1-based audit numbering) of slides that carry a baked full-bleed
    background — i.e. the bug-condition 'content-baked full-bleed' slides."""
    out = []
    for idx, slide in enumerate(prs.slides, 1):
        fbs = _fullbleed_picture_boxes(slide)
        if fbs and _picture_is_baked(slide, fbs[0]):
            out.append((idx, slide))
    return out


# ==========================================================================
# Tests — each encodes one facet of expectedBehavior (FAIL on unfixed code).
# ==========================================================================
def test_native_text_overlaps_baked_fullbleed_background():
    """(1) 네이티브 텍스트 ↔ 콘텐츠 구워진 풀블리드 배경 면적 겹침 < 10%.
    EXPECTED(post-fix). 미수정 코드: 본문 TEXT_BOX(0.6,1.6,12.1,5.4)가 구워진
    배경 위 ~100% 겹침 → FAIL (Req 1.1)."""
    pptx, _ = _generate_deck()
    prs = Presentation(pptx)

    baked = _content_slides_with_baked_fullbleed(prs)
    # Setup sanity (independent of the fix): the audited dual-layer path must be
    # exercised — at least one content slide has a content-baked full-bleed bg.
    assert baked, ("전제 실패 — 콘텐츠가 구워진 풀블리드 배경 슬라이드가 없음"
                   "(HTML 베이크 경로 미구동).")

    worst = 0.0
    worst_detail = None
    worst_idx = None
    for idx, slide in baked:
        r, detail = _worst_text_over_fullbleed(slide)
        if r > worst:
            worst, worst_detail, worst_idx = r, detail, idx

    assert worst < THRESHOLD, (
        "결함 1.1 재현 — 네이티브 텍스트가 콘텐츠 구워진 풀블리드 배경 위에 겹친다 "
        f"(슬라이드 {worst_idx}, 겹침 {worst * 100:.0f}% ≥ 임계 {THRESHOLD * 100:.0f}%).\n"
        f"  {worst_detail}"
    )


def test_title_appears_exactly_once_per_slide():
    """(2) 제목 1회 표시. EXPECTED(post-fix). 미수정 코드: 구워진 배경 제목 +
    네이티브 제목 TEXT_BOX(0.6,0.3,12.1,1.0) = 2회 → FAIL (Req 1.2)."""
    pptx, _ = _generate_deck()
    prs = Presentation(pptx)

    baked = _content_slides_with_baked_fullbleed(prs)
    assert baked, "전제 실패 — 구워진 풀블리드 배경 슬라이드가 없음."

    # Map deck title → slide by matching the native title shape text.
    offenders = []
    titles = [s["title"] for s in _CONTENT_SLIDES]
    for idx, slide in baked:
        for title in titles:
            count, native, bk = _title_occurrences(slide, title)
            if native >= 1:  # this slide carries that title natively
                if count != 1:
                    offenders.append((idx, title, count, native, bk))
                break

    assert not offenders, (
        "결함 1.2 재현 — 제목이 한 번보다 많이 표시된다(구워진 배경 제목 + 네이티브 제목).\n"
        + "\n".join(
            f"  슬라이드 {i}: '{t}' → {c}회(native {n} + baked {b})"
            for (i, t, c, n, b) in offenders)
    )


def test_all_shapes_within_slide_bounds():
    """(3) 모든 도형이 슬라이드 경계(0,0,13.333,7.5) 안. EXPECTED(post-fix).
    미수정 코드: 표지 네이티브 장식 OVAL top=-1.5(off-slide) → FAIL (Req 1.4)."""
    pptx, _ = _generate_deck()
    prs = Presentation(pptx)

    offenders = []
    for idx, slide in enumerate(prs.slides, 1):
        for nm, r in _offslide_shapes(slide):
            offenders.append((idx, nm, r))

    assert not offenders, (
        "결함 1.4 재현 — 슬라이드 경계 밖(음수-top/off-slide) 도형이 존재한다.\n"
        + "\n".join(f"  슬라이드 {i}: {nm} rect(l,t,w,h)={r}" for (i, nm, r) in offenders)
    )


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
