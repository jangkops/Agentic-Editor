"""배경 이미지에 텍스트가 구워졌는지 + 오버레이 텍스트와 중복되는지 진단.

각 슬라이드의 풀블리드 PICTURE 를 추출해:
  - 이미지 크기/종횡비
  - '텍스트가 많은 렌더(HTML 딥렌더)' 신호: 고유색 수, 엣지(글자 윤곽) 밀도
  - 사진/일러스트(텍스트 없음) 신호와 구분
그리고 슬라이드의 오버레이(편집) 텍스트 내용을 함께 출력해 육안 대조한다.
"""
from __future__ import annotations

import io
import sys
import zipfile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False


def _overlay_texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(sh, "has_text_frame", False):
            t = (sh.text_frame.text or "").strip()
            if t:
                out.append(t.replace("\n", " ⏎ "))
    return out


def _img_stats(raw):
    if not _HAS_PIL:
        return None
    try:
        im = Image.open(io.BytesIO(raw)).convert("RGB")
    except Exception:
        return None
    w, h = im.size
    small = im.resize((min(w, 480), min(h, 270)))
    px = list(small.getdata())
    uniq = len(set(px))
    # 엣지 밀도: 인접 픽셀 밝기 급변 비율(글자 윤곽이 많으면↑)
    gray = [(int(r) * 299 + int(g) * 587 + int(b) * 114) // 1000 for r, g, b in px]
    sw = small.size[0]
    edges = 0
    total = 0
    for y in range(small.size[1]):
        row = y * sw
        for x in range(sw - 1):
            total += 1
            if abs(gray[row + x] - gray[row + x + 1]) > 40:
                edges += 1
    edge_pct = (edges / total * 100.0) if total else 0.0
    return {"w": w, "h": h, "uniq": uniq, "edge_pct": round(edge_pct, 1)}


def audit(path):
    prs = Presentation(path)
    zf = zipfile.ZipFile(path)
    media = {n: zf.read(n) for n in zf.namelist() if n.startswith("ppt/media/")}
    print(f"=== 이미지/텍스트 중복 점검: {path}")
    print(f"media 파일 {len(media)}개, PIL={_HAS_PIL}\n")

    # rId → media 파일 매핑을 위해 각 슬라이드의 part 사용
    for idx, slide in enumerate(prs.slides, 1):
        overlay = _overlay_texts(slide)
        print(f"[슬라이드 {idx}] 오버레이(편집) 텍스트 {len(overlay)}개:")
        for t in overlay:
            disp = t if len(t) <= 90 else t[:90] + "…"
            print(f"    · {disp}")
        # 이미지 추출
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = sh.image.blob
            except Exception:
                blob = None
            if not blob:
                continue
            stats = _img_stats(blob)
            if stats:
                # 판정: 고유색 많고 엣지 밀도 높으면 '텍스트 풍부한 렌더(HTML 딥렌더)'
                texty = stats["edge_pct"] >= 8.0 and stats["uniq"] >= 3000
                verdict = ("텍스트 구워짐 의심(HTML 딥렌더/도표)" if texty
                           else "사진·일러스트(텍스트 적음) 가능성")
                print(f"    [배경이미지] {stats['w']}x{stats['h']} "
                      f"고유색={stats['uniq']} 엣지={stats['edge_pct']}% → {verdict}")
        print()


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: audit_pptx_images.py <pptx path>")
        raise SystemExit(2)
    audit(sys.argv[1])
