"""Parity_Scorer — 디자인 밀도 패리티 기계 판정기.

렌더된 슬라이드 HTML 문자열에서 고정 Parity_Checklist 마커의 존재 수를
세어 Density_Score 를 계산하고, 카테고리별 고정 Reference_Score 와 비교해
합격/불합격과 미충족 항목을 보고한다.

순수 함수 / 네트워크 0 / 외부 의존 없음.

Spec: pptx-design-density-parity (Requirement 5)
"""

from __future__ import annotations

from typing import Dict, List, Tuple

# 표지용 Parity_Checklist (총 7항목) — (name, marker_substring)
# Requirement 5.1
COVER_CHECKLIST: List[Tuple[str, str]] = [
    ("icon_badge", 'class="cover-icon-badge"'),
    ("notice_chip", 'class="notice-chip"'),
    ("accent_head", 'class="accent-span"'),
    ("step_grid", 'class="step-card-grid"'),
    ("accent_bar", 'class="accent-bar"'),
    ("corner_glow", 'class="corner-glow"'),
    ("footer", 'class="footer"'),
]

# 본문용 Parity_Checklist (총 8항목) — (name, marker_substring)
# Requirement 5.1
BODY_CHECKLIST: List[Tuple[str, str]] = [
    ("section_header", 'class="section-header-bar"'),
    ("contact_box", 'class="contact-box"'),
    ("note_callout", 'class="note-callout"'),
    ("link_chip", 'class="link-chip"'),
    ("numbered_item", 'class="numbered-item"'),
    ("notice_tab", 'class="notice-tab"'),
    ("slide_footer", 'class="slide-footer"'),
    ("figure_slot", 'class="figure-slot"'),
]

# 고정 Reference_Score (Genspark 참조 충족 항목 수). Requirement 5.3
COVER_REFERENCE_SCORE = 6
BODY_REFERENCE_SCORE = 6

_CHECKLISTS: Dict[str, List[Tuple[str, str]]] = {
    "cover": COVER_CHECKLIST,
    "body": BODY_CHECKLIST,
}

_REFERENCE_SCORES: Dict[str, int] = {
    "cover": COVER_REFERENCE_SCORE,
    "body": BODY_REFERENCE_SCORE,
}


def score(html: str, category: str) -> dict:
    """렌더 HTML 에서 Parity_Checklist 마커를 세어 Density_Score 를 산출한다.

    Args:
        html: 렌더된 슬라이드 HTML 문자열.
        category: "cover" 또는 "body".

    Returns:
        {
          category: str,
          density_score: int,        # 0..total
          reference_score: int,      # 0..total (고정)
          total: int,
          passed: bool,              # density_score >= reference_score
          items: [{name: str, present: bool}, ...],
          missing: [name, ...],
        }

    Raises:
        ValueError: category 가 {"cover","body"} 가 아니거나,
                    html 이 None/빈 문자열인 경우 (Requirement 5.9 인접 정책).
    """
    if category not in _CHECKLISTS:
        raise ValueError(
            f"category must be one of {sorted(_CHECKLISTS)}, got {category!r}"
        )
    # Requirement 5.9 인접 정책: 입력 누락/빈 문자열 → 점수 미산출
    if html is None or html == "":
        raise ValueError("html must be a non-empty string")

    checklist = _CHECKLISTS[category]
    reference_score = _REFERENCE_SCORES[category]
    total = len(checklist)

    items: List[Dict[str, object]] = []
    missing: List[str] = []
    density_score = 0
    for name, marker in checklist:
        present = marker in html
        items.append({"name": name, "present": present})
        if present:
            density_score += 1
        else:
            missing.append(name)

    return {
        "category": category,
        "density_score": density_score,
        "reference_score": reference_score,
        "total": total,
        "passed": density_score >= reference_score,
        "items": items,
        "missing": missing,
    }


# ===========================================================================
# Native_Shape_Adapter — pptx-native-density-render (작업 5.1)
# ---------------------------------------------------------------------------
# 기존 score(html, category)는 HTML 마커 문자열을 센다. 네이티브 .pptx 슬라이드
# (python-pptx Slide 객체)를 채점하려면 셰이프 트리에서 "시각 요소 존재"를
# 검출해 체크리스트 마커로 환산하는 얇은 어댑터가 필요하다.
#
# 가산적 추가 원칙(설계 §Components §4):
#   - 기존 score()/_CHECKLISTS/_REFERENCE_SCORES/COVER_CHECKLIST/BODY_CHECKLIST는
#     일절 변경하지 않고 재사용한다(신규 토큰/체크리스트 신설 금지).
#   - 반환 dict 형식은 score()와 동일(category/density_score/reference_score/
#     total/passed/items/missing).
#   - 빈 입력(slide=None) 또는 미지원 카테고리(cover·body 외)는 ValueError
#     (기존 score()의 5.5 정책과 일관).
#
# 매핑 방식(설계 매핑 표 정합): native_layout_renderer는 PlacedShape.role을
# 부여하지만 그 role 메타는 .pptx 파일에 보존되지 않으므로, 채점은 python-pptx
# 셰이프의 shape_type/auto_shape_type/텍스트/색/기하 특징으로 각 체크리스트 항목에
# 대응하는 시각 요소 존재를 휴리스틱 검출한다. 정확한 1:1이 어려운 항목은 가장
# 근접한 시각 요소 존재로 합리적으로 매핑한다(아래 각 검출기 주석 참조):
#   cover : icon_badge / notice_chip / accent_head / step_grid /
#           accent_bar / corner_glow / footer
#   body  : section_header / contact_box / note_callout / link_chip /
#           numbered_item / notice_tab / slide_footer / figure_slot
# ===========================================================================

_EMU_PER_INCH = 914400.0
# 슬라이드 경계(16:9, 인치) — layout_geometry.SLIDE_RECT 와 정합.
_NS_SLIDE_W_IN = 13.333
_NS_SLIDE_H_IN = 7.5


def _ns_emu_to_in(v) -> float:
    """EMU(또는 Emu 객체) → 인치. 변환 불가 시 0.0."""
    try:
        return float(int(v)) / _EMU_PER_INCH
    except Exception:
        return 0.0


def _ns_rect(shp):
    """셰이프의 (left, top, width, height)를 인치 튜플로. 산출 불가 시 None."""
    try:
        l = _ns_emu_to_in(shp.left)
        t = _ns_emu_to_in(shp.top)
        w = _ns_emu_to_in(shp.width)
        h = _ns_emu_to_in(shp.height)
    except Exception:
        return None
    if w <= 0 or h <= 0:
        return None
    return (l, t, w, h)


def _ns_fill_rgb(shp):
    """솔리드 채움 색을 (r, g, b)로. 솔리드 RGB가 아니면(테마/그라데이션/없음) None."""
    try:
        from pptx.enum.dml import MSO_FILL_TYPE  # type: ignore
    except Exception:
        MSO_FILL_TYPE = None  # type: ignore
    try:
        fill = shp.fill
        if MSO_FILL_TYPE is not None and fill.type != MSO_FILL_TYPE.SOLID:
            return None
        rgb = fill.fore_color.rgb
        return (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    except Exception:
        return None


def _ns_has_line(shp) -> bool:
    """가시 테두리(line)가 설정됐는지(폭 또는 색). 판정 불가 시 False."""
    try:
        ln = shp.line
        if ln.width and int(ln.width) > 0:
            return True
    except Exception:
        pass
    try:
        _ = shp.line.color.rgb  # RGB 색이 설정돼 있으면 접근 성공.
        return True
    except Exception:
        return False


def _ns_auto_name(shp):
    """auto_shape_type 이름 문자열(대문자). 오토셰이프가 아니면 None.

    예: 'ROUNDED_RECTANGLE', 'RECTANGLE', 'OVAL'. ('ROUNDED_RECTANGLE'은
    'RECTANGLE'을 부분문자열로 포함하므로 검출 시 ROUNDED를 먼저 확인할 것.)
    """
    try:
        at = shp.auto_shape_type
    except Exception:
        return None
    if at is None:
        return None
    name = getattr(at, "name", None)
    if name:
        return str(name).upper()
    s = str(at).upper()
    # 'ROUNDED_RECTANGLE (5)' 형태 → 토큰만.
    return s.split("(")[0].strip().replace(" ", "_")


def _ns_text(shp) -> str:
    """셰이프 텍스트(프레임/표 셀 포함). 없으면 ''. 예외 시 ''."""
    out = []
    try:
        if shp.has_text_frame:
            out.append(shp.text_frame.text or "")
    except Exception:
        pass
    try:
        if getattr(shp, "has_table", False):
            for row in shp.table.rows:
                for cell in row.cells:
                    out.append(cell.text or "")
    except Exception:
        pass
    return "\n".join(s for s in out if s).strip()


def _ns_max_font_pt(shp):
    """텍스트 프레임 내 최대 run 폰트 크기(pt). 없으면 None."""
    try:
        if not shp.has_text_frame:
            return None
    except Exception:
        return None
    mx = None
    try:
        for p in shp.text_frame.paragraphs:
            for r in p.runs:
                try:
                    sz = r.font.size
                    if sz is not None:
                        pt = float(sz.pt)
                        mx = pt if mx is None else max(mx, pt)
                except Exception:
                    continue
    except Exception:
        return mx
    return mx


def _ns_shape_kind(shp) -> str:
    """셰이프 대분류: 'picture'|'autoshape'|'textbox'|'table'|'placeholder'|'other'."""
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        MSO_SHAPE_TYPE = None  # type: ignore
    try:
        if getattr(shp, "has_table", False):
            return "table"
    except Exception:
        pass
    st = None
    try:
        st = shp.shape_type
    except Exception:
        st = None
    if MSO_SHAPE_TYPE is not None and st is not None:
        try:
            if st == MSO_SHAPE_TYPE.PICTURE:
                return "picture"
            if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                return "autoshape"
            if st == MSO_SHAPE_TYPE.TEXT_BOX:
                return "textbox"
            if st == MSO_SHAPE_TYPE.TABLE:
                return "table"
            if st == MSO_SHAPE_TYPE.PLACEHOLDER:
                return "placeholder"
        except Exception:
            pass
    # 폴백: 오토셰이프 여부를 auto_shape_type 접근 성공으로 추정.
    if _ns_auto_name(shp) is not None:
        return "autoshape"
    try:
        if shp.has_text_frame:
            return "textbox"
    except Exception:
        pass
    return "other"


def _ns_collect(slide) -> list:
    """슬라이드 셰이프 트리를 채점용 정규화 dict 리스트로 수집(순수 데이터).

    각 항목: {kind, auto, text, has_text, rect, fill, line, max_pt}.
    이후 검출기는 이 dict 리스트만으로 동작하여 python-pptx 의존을 격리한다.
    """
    infos: list = []
    try:
        shapes = list(slide.shapes)
    except Exception:
        return infos
    for shp in shapes:
        try:
            text = _ns_text(shp)
            info = {
                "kind": _ns_shape_kind(shp),
                "auto": _ns_auto_name(shp),
                "text": text,
                "has_text": bool(text),
                "rect": _ns_rect(shp),
                "fill": _ns_fill_rgb(shp),
                "line": _ns_has_line(shp),
                "max_pt": _ns_max_font_pt(shp),
            }
        except Exception:
            continue
        infos.append(info)
    return infos


# --- 색 판정 헬퍼 (순수, RGB 튜플 입력) --------------------------------------


def _ns_luma(rgb) -> float:
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _ns_is_dark(rgb) -> bool:
    """어두운 채움(섹션 헤더 다크 막대 등). 휘도 < 110."""
    return rgb is not None and _ns_luma(rgb) < 110.0


def _ns_is_near_white(rgb) -> bool:
    return rgb is not None and all(c >= 235 for c in rgb)


def _ns_is_near_black(rgb) -> bool:
    return rgb is not None and all(c <= 28 for c in rgb)


def _ns_is_chromatic(rgb) -> bool:
    """유채색(accent/primary 계열) — 검정/흰색/회색이 아닌 색. 채널 분산 > 24."""
    if rgb is None:
        return False
    if _ns_is_near_white(rgb) or _ns_is_near_black(rgb):
        return False
    return (max(rgb) - min(rgb)) > 24


def _ns_is_light_gray(rgb) -> bool:
    """그림 슬롯 플레이스홀더(중립 회색, 예: F0F2F5=(240,242,245))용 — 밝고 무채색.

    네이티브 emit_figure_slot 플레이스홀더 색을 검출 대상에 포함하기 위해 상한을
    250 으로 두어(순백 255 만 제외) 밝은 중립 회색을 인정한다."""
    if rgb is None:
        return False
    return (max(rgb) - min(rgb)) <= 18 and 195 <= min(rgb) <= 250


# --- cover 카테고리 검출기 (infos 리스트 입력) -------------------------------


def _ns_cover_icon_badge(infos) -> bool:
    """아이콘 배지 ← 작고 정사각형에 가까운 배지형 도형(OVAL/ROUNDED_RECTANGLE),
    유채색 채움 또는 텍스트 보유, 최대변 ≤ 1.3in."""
    for i in infos:
        auto = i["auto"]
        rect = i["rect"]
        if i["kind"] == "autoshape" and auto and rect and (
            "OVAL" in auto or "ROUNDED_RECTANGLE" in auto
        ):
            l, t, w, h = rect
            ar = w / h if h else 0.0
            if 0.5 <= ar <= 2.0 and max(w, h) <= 1.3 and (
                i["has_text"] or _ns_is_chromatic(i["fill"])
            ):
                return True
    return False


def _ns_cover_notice_chip(infos) -> bool:
    """노티스 칩 ← 작은 pill형 둥근 사각형(eyebrow pill 등), 높이 ≤ 0.7in,
    폭 ≥ 1.0in, 채움 또는 텍스트 보유."""
    for i in infos:
        auto = i["auto"]
        rect = i["rect"]
        if auto and "ROUNDED_RECTANGLE" in auto and rect:
            l, t, w, h = rect
            if h <= 0.7 and w >= 1.0 and (i["has_text"] or i["fill"] is not None):
                return True
    return False


def _ns_cover_accent_head(infos) -> bool:
    """강조 헤드 ← 대형 제목 텍스트(최대 폰트 ≥ 26pt를 가진 텍스트 셰이프)."""
    for i in infos:
        if i["has_text"] and i["max_pt"] is not None and i["max_pt"] >= 26.0:
            return True
    return False


def _ns_cover_step_grid(infos) -> bool:
    """스텝/카드 그리드 ← 일정 크기 이상의 둥근 사각형 카드가 2개 이상."""
    cnt = 0
    for i in infos:
        auto = i["auto"]
        rect = i["rect"]
        if auto and "ROUNDED_RECTANGLE" in auto and rect:
            l, t, w, h = rect
            if w >= 1.2 and h >= 0.8:
                cnt += 1
    return cnt >= 2


def _ns_cover_accent_bar(infos) -> bool:
    """액센트 바 ← 얇고 긴 막대: 세로(w ≤ 0.6, h ≥ 2.0) 또는 가로(h ≤ 0.3, w ≥ 1.5)."""
    for i in infos:
        rect = i["rect"]
        if i["kind"] == "autoshape" and rect:
            l, t, w, h = rect
            if (w <= 0.6 and h >= 2.0) or (h <= 0.3 and w >= 1.5):
                return True
    return False


def _ns_cover_corner_glow(infos) -> bool:
    """코너 글로우 ← 대형 장식 원(OVAL, 최대변 ≥ 1.5in)."""
    for i in infos:
        auto = i["auto"]
        rect = i["rect"]
        if auto and "OVAL" in auto and rect:
            l, t, w, h = rect
            if max(w, h) >= 1.5:
                return True
    return False


def _ns_cover_footer(infos) -> bool:
    """푸터 ← 슬라이드 하단부(top ≥ slideH - 1.2in)에 위치한 셰이프 존재."""
    for i in infos:
        rect = i["rect"]
        if rect and rect[1] >= (_NS_SLIDE_H_IN - 1.2):
            return True
    return False


# --- body 카테고리 검출기 (infos 리스트 입력) --------------------------------


def _ns_body_section_header(infos) -> bool:
    """섹션 헤더 ← 어두운 채움의 넓은 막대(autoshape, dark, w ≥ 2.5, w ≥ h*2.2, h ≤ 1.4).
    (native emit_section_header_bar 의 다크 ROUNDED_RECTANGLE 막대.)"""
    for i in infos:
        rect = i["rect"]
        if i["kind"] == "autoshape" and rect and _ns_is_dark(i["fill"]):
            l, t, w, h = rect
            if w >= 2.5 and h <= 1.4 and w >= h * 2.2:
                return True
    return False


def _ns_body_contact_box(infos) -> bool:
    """연락처 박스 ← 테두리 있는 둥근 사각형 박스 + 좌측 얇은 세로 액센트 바 공존.
    (native emit_contact_box 의 틴트 박스 + 세로 accent 바 구조.)"""
    has_thin_vbar = any(
        i["kind"] == "autoshape" and i["rect"]
        and i["rect"][2] <= 0.25 and i["rect"][3] >= 0.6
        and _ns_is_chromatic(i["fill"])
        for i in infos
    )
    has_bordered_box = any(
        i["auto"] and "ROUNDED_RECTANGLE" in i["auto"] and i["line"] and i["rect"]
        and i["rect"][2] >= 2.0 and 0.5 <= i["rect"][3] <= 3.5
        for i in infos
    )
    return has_thin_vbar and has_bordered_box


def _ns_body_note_callout(infos) -> bool:
    """노트 콜아웃 ← 'NOTICE' 라벨 텍스트 존재(native emit_note_callout note_label).
    폴백: accent 테두리를 가진 틴트 둥근 박스."""
    for i in infos:
        if i["has_text"] and "NOTICE" in i["text"].upper():
            return True
    # 폴백: accent 유채색 테두리 + 둥근 박스(라벨 텍스트가 비어도 구조로 인정).
    for i in infos:
        if i["auto"] and "ROUNDED_RECTANGLE" in i["auto"] and i["line"] and i["rect"]:
            l, t, w, h = i["rect"]
            if w >= 2.0 and 0.4 <= h <= 3.0:
                return True
    return False


def _ns_body_link_chip(infos) -> bool:
    """링크 칩 ← 텍스트를 가진 작은 pill형 둥근 사각형(h ≤ 0.55, 0.6 ≤ w ≤ 4.0)."""
    for i in infos:
        auto = i["auto"]
        rect = i["rect"]
        if auto and "ROUNDED_RECTANGLE" in auto and rect and i["has_text"]:
            l, t, w, h = rect
            if h <= 0.55 and 0.6 <= w <= 4.0:
                return True
    return False


def _ns_body_numbered_item(infos) -> bool:
    """번호 항목 ← 숫자 텍스트를 가진 원형 배지(OVAL + 숫자).
    (native emit_numbered_list / section_header 배지의 번호 oval.)"""
    for i in infos:
        auto = i["auto"]
        if auto and "OVAL" in auto and i["has_text"] and i["text"].strip().isdigit():
            return True
    return False


def _ns_body_notice_tab(infos) -> bool:
    """노티스 탭 ← 작은 유채색(accent) 솔리드 탭/마커(면적 ≤ 2.0, w ≤ 2.5, h ≤ 0.7)."""
    for i in infos:
        rect = i["rect"]
        if i["kind"] == "autoshape" and rect and _ns_is_chromatic(i["fill"]):
            l, t, w, h = rect
            if w <= 2.5 and h <= 0.7 and (w * h) <= 2.0:
                return True
    return False


def _ns_body_slide_footer(infos) -> bool:
    """슬라이드 푸터 ← 하단부(top ≥ slideH - 1.0in)에 위치한 작은 텍스트 셰이프(h ≤ 0.7)."""
    for i in infos:
        rect = i["rect"]
        if i["has_text"] and rect:
            l, t, w, h = rect
            if t >= (_NS_SLIDE_H_IN - 1.0) and h <= 0.7:
                return True
    return False


def _ns_body_figure_slot(infos) -> bool:
    """그림 슬롯 ← 그림(picture) 또는 텍스트 없는 중립 회색 사각형 플레이스홀더
    (native emit_figure_slot: 손실-0 이미지 또는 회색 RECTANGLE)."""
    for i in infos:
        if i["kind"] == "picture":
            return True
        auto = i["auto"]
        rect = i["rect"]
        if (
            i["kind"] == "autoshape" and auto and "RECTANGLE" in auto
            and "ROUNDED_RECTANGLE" not in auto and not i["has_text"] and rect
            and _ns_is_light_gray(i["fill"])
        ):
            l, t, w, h = rect
            if w >= 1.5 and h >= 1.0:
                return True
    return False


_NATIVE_DETECTORS = {
    "cover": {
        "icon_badge": _ns_cover_icon_badge,
        "notice_chip": _ns_cover_notice_chip,
        "accent_head": _ns_cover_accent_head,
        "step_grid": _ns_cover_step_grid,
        "accent_bar": _ns_cover_accent_bar,
        "corner_glow": _ns_cover_corner_glow,
        "footer": _ns_cover_footer,
    },
    "body": {
        "section_header": _ns_body_section_header,
        "contact_box": _ns_body_contact_box,
        "note_callout": _ns_body_note_callout,
        "link_chip": _ns_body_link_chip,
        "numbered_item": _ns_body_numbered_item,
        "notice_tab": _ns_body_notice_tab,
        "slide_footer": _ns_body_slide_footer,
        "figure_slot": _ns_body_figure_slot,
    },
}


def score_native_slide(slide, category: str) -> dict:
    """네이티브 .pptx 슬라이드(python-pptx Slide)를 채점한다.

    셰이프 트리에서 각 Parity_Checklist 항목에 대응하는 "시각 요소 존재"를
    휴리스틱 검출(shape_type/auto_shape_type/텍스트/색/기하)하여 Density_Score 를
    산출하고, 기존 score() 와 **동일한 dict 형식**으로 반환한다.

    기존 _CHECKLISTS / _REFERENCE_SCORES(cover=6, body=6)를 재사용하며(신규 토큰/
    체크리스트 신설 없음), score() 및 관련 상수는 변경하지 않는다(가산적 추가).

    Args:
        slide: python-pptx Slide 객체.
        category: "cover" 또는 "body".

    Returns:
        {
          category, density_score, reference_score, total,
          passed,            # density_score >= reference_score
          items: [{name, present}, ...],
          missing: [name, ...],
        }
        — score() 와 동일 키 집합.

    Raises:
        ValueError: category 가 {"cover","body"} 가 아니거나 slide 가 None 인 경우
                    (기존 score() 의 5.5 정책과 일관).
    """
    if category not in _CHECKLISTS:
        raise ValueError(
            f"category must be one of {sorted(_CHECKLISTS)}, got {category!r}"
        )
    # Req 5.5 인접 정책: 빈 입력(None) → 점수 미산출.
    if slide is None:
        raise ValueError("slide must be a non-None python-pptx Slide")

    detectors = _NATIVE_DETECTORS[category]
    checklist = _CHECKLISTS[category]
    reference_score = _REFERENCE_SCORES[category]
    total = len(checklist)

    infos = _ns_collect(slide)

    items: List[Dict[str, object]] = []
    missing: List[str] = []
    density_score = 0
    for name, _marker in checklist:
        det = detectors.get(name)
        present = bool(det(infos)) if det is not None else False
        items.append({"name": name, "present": present})
        if present:
            density_score += 1
        else:
            missing.append(name)

    return {
        "category": category,
        "density_score": density_score,
        "reference_score": reference_score,
        "total": total,
        "passed": density_score >= reference_score,
        "items": items,
        "missing": missing,
    }
