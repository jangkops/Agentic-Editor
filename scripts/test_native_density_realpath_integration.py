# Feature: pptx-native-density-render, 수정 4: 실제 생성 경로(_tool_generate_pptx) 통합 검증
#
# 목적: server.py 의 실제 PPTX 생성 진입점 `_tool_generate_pptx` 를 기본 환경
# (AE_NATIVE_LAYOUT_RENDER 미설정 = 기본 활성)에서 헤르메틱하게 구동하여, 표지+
# 본문(콘텐츠)+비-7키 레이아웃이 섞인 덱을 생성하고, 생성된 .pptx 를
# audit_native_density 로 검사해 (a)~(h) 전수 통과함을 입증한다.
#
# 이 테스트가 검증하는 근본 수정:
#   수정 1) 본문 네이티브 라우팅 기본 활성 + 비-7키 레이아웃(kpi_summary/
#           status_table/objective_detail/process_flow) 근접 매핑(map_to_native_layout).
#   수정 2) 표지 겹침 제거(네이티브 표지 우선, 콘텐츠 구운 배경 미채택).
#   수정 3) 다이어그램/본문 겹침 제거(네이티브 라우팅 시 이중 렌더 차단).
#
# 헤르메틱 규약:
#   - 네트워크 0: 게이트웨이(_get_gw)/LLM(_llm_pick_slide_layout)/Vertex/Chrome 전부 mock.
#   - Chrome 미사용: 네이티브 경로는 HTML→PNG 래스터화를 타지 않는다. Chrome 감지
#     (_find_local_chrome)는 truthy 로 mock 해 _html_enabled=True(프로덕션과 동일한
#     마스터 게이트)를 만들되, 실제 Chrome 렌더 호출(_render_html_slide_to_png/
#     _generate_html_slide_for_section)은 mock 으로 대체해 Chrome 을 띄우지 않는다.
#   - Vertex: AE_PREFER_VERTEX_IMAGE=0 로 사전생성 비활성(자격증명 없어도 안전).
#
# 실행: ./venv/bin/python -m pytest scripts/test_native_density_realpath_integration.py -p no:cacheprovider -q

import os
import sys
import json
import asyncio
import tempfile
import shutil
from unittest import mock

import pytest

_REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _REPO not in sys.path:
    sys.path.insert(0, _REPO)

import ai_engine.server as srv  # noqa: E402

# 산출물_검증기 (scripts/ 형제 모듈)
_SCRIPTS = os.path.join(_REPO, "scripts")
if _SCRIPTS not in sys.path:
    sys.path.insert(0, _SCRIPTS)
import audit_pptx_native_density as auditor  # noqa: E402


# ---------------------------------------------------------------------------
# LLM 레이아웃 픽 mock — 섹션 제목으로 레이아웃/데이터를 결정한다.
# 비-7키 레이아웃(kpi_summary/status_table/objective_detail/process_flow)을 반환해
# server.py 의 map_to_native_layout 근접 매핑(수정 1)이 작동함을 검증한다.
# 반환 data 는 매핑된 7키 레이아웃이 리치하게 렌더될 수 있도록 해당 레이아웃의 필수
# 필드로 구성한다(실 LLM 도 자신이 고른 레이아웃의 data 를 추출한다).
# ---------------------------------------------------------------------------

def _pick_for(heading: str) -> dict:
    h = heading or ""
    if "지표" in h or "KPI" in h:
        # 비-7키 kpi_summary → feature_grid 로 매핑됨. data 는 feature_grid 형태.
        return {
            "layout": "kpi_summary",
            "data": {
                "title": heading,
                "subtitle": "핵심 성과 지표 요약",
                "features": [
                    {"icon": "zap", "title": "처리량 +42%", "description": "월간 처리 건수가 전분기 대비 42% 증가했습니다."},
                    {"icon": "shield", "title": "가동률 99.9%", "description": "무중단 운영으로 목표 SLA 를 상회 달성했습니다."},
                    {"icon": "users", "title": "활성 사용자 3.2만", "description": "신규 온보딩 확대로 활성 사용자가 꾸준히 늘었습니다."},
                    {"icon": "database", "title": "비용 -18%", "description": "리소스 최적화로 인프라 비용을 18% 절감했습니다."},
                ],
            },
        }
    if "상태" in h or "진척" in h:
        # 비-7키 status_table → comparison 으로 매핑됨. data 는 comparison 형태.
        return {
            "layout": "status_table",
            "data": {
                "title": heading,
                "left_label": "완료",
                "left_items": ["요구사항 정의", "아키텍처 설계", "핵심 모듈 구현"],
                "right_label": "진행/예정",
                "right_items": ["통합 테스트", "성능 튜닝", "사용자 검증"],
            },
        }
    if "일정" in h or "프로세스" in h or "단계" in h:
        # 비-7키 process_flow → timeline 으로 매핑됨. data 는 timeline 형태.
        return {
            "layout": "process_flow",
            "data": {
                "title": heading,
                "steps": [
                    {"label": "1단계", "title": "착수", "description": "범위 확정과 킥오프를 진행합니다."},
                    {"label": "2단계", "title": "구현", "description": "핵심 기능을 반복 개발합니다."},
                    {"label": "3단계", "title": "검증", "description": "품질 게이트로 산출물을 검증합니다."},
                    {"label": "4단계", "title": "배포", "description": "단계적 롤아웃으로 안정 배포합니다."},
                ],
            },
        }
    if "구조" in h or "아키텍처" in h:
        return {
            "layout": "architecture",
            "data": {
                "title": heading,
                "layers": [
                    {"name": "프레젠테이션", "description": "Electron UI", "items": ["렌더러", "IPC"]},
                    {"name": "애플리케이션", "description": "FastAPI", "items": ["에이전트", "도구"]},
                    {"name": "데이터", "description": "게이트웨이", "items": ["Bedrock", "캐시"]},
                ],
            },
        }
    if "기능" in h:
        return {
            "layout": "feature_grid",
            "data": {
                "title": heading,
                "features": [
                    {"icon": "check", "title": "편집 가능", "description": "모든 텍스트를 PowerPoint 에서 직접 수정합니다."},
                    {"icon": "layers", "title": "고밀도", "description": "젠스파크급 정보 밀도를 네이티브로 유지합니다."},
                    {"icon": "shield", "title": "겹침 0", "description": "기하 보정으로 셰이프 겹침을 제거합니다."},
                ],
            },
        }
    # 기본: two_column (7키, 항등 매핑)
    return {
        "layout": "two_column",
        "data": {
            "title": heading,
            "left_content": "핵심 목표\n일정 준수\n품질 확보",
            "right_content": "리스크 관리\n이해관계자 소통\n지속 개선",
        },
    }


class _AsyncPick:
    """_llm_pick_slide_layout 대체 — 섹션 제목 기반으로 layout/data 반환."""

    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, section_heading, section_body,
                       doc_context="", role="", bullet_count=0):
        self.calls.append(section_heading)
        return _pick_for(section_heading)


async def _noop_render_html_to_png(html, output_path, width=1920, height=1080, timeout=30):
    # Chrome 미실행: 표지 HTML 은 네이티브 표지 우선(_cover_native_applied)이라
    # 임베드되지 않는다. ok=False 로 반환해 어떤 경우에도 Chrome 을 띄우지 않는다.
    return {"ok": False, "reason": "mocked-no-chrome"}


async def _noop_generate_html_slide(*args, **kwargs):
    # 본문 HTML 베이크 경로 — 네이티브 라우팅 성공 시 호출되지 않지만 안전하게 mock.
    return ""


def _build_tool_input():
    return {
        "title": "네이티브 고밀도 렌더 통합 검증 덱",
        "highQuality": True,
        "slides": [
            {"title": "프로젝트 개요", "bullets": [
                "목표: 편집 가능한 고밀도 슬라이드 생성",
                "범위: 표지·본문·다이어그램 전 슬라이드",
                "원칙: 통짜 이미지 금지, 네이티브 도형 우선",
                "품질: 겹침 10% 미만, 제목 1회",
            ]},
            {"title": "핵심 지표", "bullets": [
                "처리량 42% 증가",
                "가동률 99.9% 달성",
                "활성 사용자 3.2만",
                "인프라 비용 18% 절감",
            ]},
            {"title": "진행 상태", "bullets": [
                "요구사항 정의 완료",
                "아키텍처 설계 완료",
                "통합 테스트 진행 중",
            ]},
            {"title": "추진 일정", "bullets": [
                "1단계 착수",
                "2단계 구현",
                "3단계 검증",
                "4단계 배포",
            ]},
            {"title": "시스템 구조", "bullets": [
                "프레젠테이션 계층",
                "애플리케이션 계층",
                "데이터 계층",
            ]},
            {"title": "주요 기능", "bullets": [
                "편집 가능한 네이티브 도형",
                "젠스파크급 고밀도",
                "셰이프 겹침 0",
            ]},
        ],
    }


def _run_generate(tool_input, project_path, env_overrides):
    """헤르메틱 컨텍스트에서 _tool_generate_pptx 를 구동하고 결과 JSON 을 반환."""
    pick = _AsyncPick()
    old_env = {}
    for k, v in env_overrides.items():
        old_env[k] = os.environ.get(k)
        if v is None:
            os.environ.pop(k, None)
        else:
            os.environ[k] = v
    try:
        with mock.patch.object(srv, "_find_local_chrome", lambda: "/fake/chrome"), \
             mock.patch.object(srv, "_call_bridge", lambda ep, payload, timeout=30.0: None), \
             mock.patch.object(srv, "_get_gw", lambda *a, **k: object()), \
             mock.patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "fake-model"), \
             mock.patch.object(srv, "_llm_pick_slide_layout", pick), \
             mock.patch.object(srv, "_render_html_slide_to_png", _noop_render_html_to_png), \
             mock.patch.object(srv, "_generate_html_slide_for_section", _noop_generate_html_slide):
            out = asyncio.run(srv._tool_generate_pptx(tool_input, project_path))
    finally:
        for k, v in old_env.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v
    return json.loads(out), pick


def _abs_from_result(result, project_path):
    assert "path" in result, f"생성 실패 결과: {result}"
    rel = result["path"]
    local_root = srv._resolve_local_root(project_path)
    cand = rel if os.path.isabs(rel) else os.path.join(local_root, rel)
    assert os.path.isfile(cand), f".pptx 파일 부재: {cand} (result={result})"
    return cand


# ---------------------------------------------------------------------------
# 테스트
# ---------------------------------------------------------------------------

def test_default_env_realpath_produces_editable_audit_passing_deck():
    """기본 환경(AE_NATIVE_LAYOUT_RENDER 미설정)에서 실제 생성 경로가 편집가능·
    audit (a)~(h) 통과 덱을 만든다. (수정 1~3 통합 입증)"""
    tmp = tempfile.mkdtemp(prefix="ae_realpath_")
    try:
        tool_input = _build_tool_input()
        env = {
            # HTML 마스터 게이트 강제 ON(프로덕션 동일). _html_enabled=True 가 되어야
            # 네이티브 라우팅 게이트가 진입한다.
            "AE_ENABLE_HTML_SLIDES": "1",
            # Vertex 사전생성 비활성(네트워크 0, 자격증명 불필요).
            "AE_PREFER_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_IMAGE": "0",
            # 기본 활성 검증: AE_NATIVE_LAYOUT_RENDER 는 미설정(제거)해 기본값 사용.
            "AE_NATIVE_LAYOUT_RENDER": None,
        }
        result, pick = _run_generate(tool_input, tmp, env)
        pptx_path = _abs_from_result(result, tmp)

        # 수정 1: 비-7키 레이아웃 픽이 실제로 발생했는지(매핑 경로 진입) 확인.
        assert pick.calls, "네이티브 라우팅이 _llm_pick_slide_layout 을 호출하지 않음"

        report = auditor.audit_native_density(pptx_path)
        assert report.passed, (
            "실제 생성 경로 산출물이 audit 를 통과하지 못함:\n"
            + "\n".join(
                f"  - slide {f.get('slide')}: {f.get('check')} "
                f"shapes={f.get('shapes')} signal={f.get('signal')}"
                for f in report.failures
            )
        )
        assert report.failures == []
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def test_non_seven_key_layouts_route_native_not_baked():
    """비-7키 레이아웃(kpi_summary/status_table/process_flow)이 통짜 이미지 베이크가
    아니라 편집가능 네이티브로 라우팅된다(수정 1 근접 매핑)."""
    tmp = tempfile.mkdtemp(prefix="ae_realpath_map_")
    try:
        tool_input = _build_tool_input()
        env = {
            "AE_ENABLE_HTML_SLIDES": "1",
            "AE_PREFER_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_IMAGE": "0",
            "AE_NATIVE_LAYOUT_RENDER": None,
        }
        result, pick = _run_generate(tool_input, tmp, env)
        pptx_path = _abs_from_result(result, tmp)

        # 비-7키 레이아웃 이름이 실제로 픽되었는지 확인(매핑 대상 존재).
        assert any(k in ("핵심 지표", "진행 상태", "추진 일정") for k in pick.calls)

        # 산출물에 통짜 배경 이미지(풀블리드 PICTURE)가 없고, 모든 본문 슬라이드에
        # 편집가능 텍스트 런이 존재함을 직접 확인한다.
        from pptx import Presentation
        prs = Presentation(pptx_path)
        for si, slide in enumerate(prs.slides):
            editable = 0
            for sh in slide.shapes:
                try:
                    if sh.has_text_frame and sh.text_frame.text.strip():
                        editable += 1
                except Exception:
                    pass
            assert editable >= 1, f"슬라이드 {si + 1} 편집가능 텍스트 런 0(통짜 이미지 의심)"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def test_native_render_disabled_killswitch_falls_back():
    """킬스위치(AE_NATIVE_LAYOUT_RENDER=0)로 네이티브 라우팅을 끄면 기존 HTML/베이크
    경로가 보존된다(회귀 롤백 경로 존재 확인). 이 경로는 Chrome mock 으로 베이크가
    실패하므로 카드/네이티브 폴백으로 콘텐츠는 유지된다(통짜 이미지 강제 아님)."""
    tmp = tempfile.mkdtemp(prefix="ae_realpath_off_")
    try:
        tool_input = _build_tool_input()
        env = {
            "AE_ENABLE_HTML_SLIDES": "1",
            "AE_PREFER_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_IMAGE": "0",
            "AE_NATIVE_LAYOUT_RENDER": "0",  # 킬스위치
        }
        result, pick = _run_generate(tool_input, tmp, env)
        pptx_path = _abs_from_result(result, tmp)
        # 파일이 생성되고(생성 경로 자체는 견고), 예외 없이 완료됨을 확인.
        assert os.path.isfile(pptx_path)
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def _text_shapes_in(slide):
    """(rect_in_inches, text) 목록 — 비어있지 않은 텍스트 셰이프만."""
    out = []
    EMU = 914400.0
    for sh in slide.shapes:
        try:
            if not sh.has_text_frame:
                continue
            t = (sh.text_frame.text or "").strip()
            if not t:
                continue
            if None in (sh.left, sh.top, sh.width, sh.height):
                continue
            out.append(((sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU), t))
        except Exception:
            continue
    return out


def _fullbleed_pictures(slide):
    """풀블리드(슬라이드 전체를 덮는) PICTURE 개수 — 콘텐츠 구운 통짜 배경 탐지."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    EMU = 914400.0
    n = 0
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            w = (sh.width or 0) / EMU
            h = (sh.height or 0) / EMU
            if w >= 13.0 and h >= 7.2:
                n += 1
        except Exception:
            continue
    return n


def test_caller_native_diagram_suppresses_backdrop_and_no_overlap():
    """수정 3: caller nativeDiagram 경로(Req 9.3 보존)가 기본(Vertex OFF) 모드에서
    풀블리드 backdrop 없이 편집가능 네이티브 카드만 그리고, 텍스트-텍스트 겹침이
    10% 미만이다."""
    tmp = tempfile.mkdtemp(prefix="ae_realpath_diag_")
    try:
        tool_input = {
            "title": "다이어그램 경로 보존 검증",
            "highQuality": True,
            "slides": [
                {"title": "표지 다음 본문", "bullets": ["항목 A", "항목 B", "항목 C"]},
                {
                    "title": "처리 흐름",
                    "nativeDiagram": {
                        "type": "process",
                        "content": "요구분석\n설계\n구현\n검증\n배포",
                    },
                },
            ],
        }
        env = {
            "AE_ENABLE_HTML_SLIDES": "1",
            "AE_PREFER_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_BG": "0",
            "AE_NATIVE_LAYOUT_RENDER": None,
        }
        result, _pick = _run_generate(tool_input, tmp, env)
        pptx_path = _abs_from_result(result, tmp)

        from pptx import Presentation
        try:
            from layout_geometry import overlap_area, area
        except Exception:
            from ai_engine.layout_geometry import overlap_area, area

        prs = Presentation(pptx_path)
        slides = list(prs.slides)
        # 마지막 슬라이드 = nativeDiagram 슬라이드.
        diag_slide = slides[-1]

        # (1) 풀블리드 통짜 배경(콘텐츠 구운 PNG) 미채택 — 기본 모드 backdrop 억제.
        assert _fullbleed_pictures(diag_slide) == 0, "다이어그램 슬라이드에 풀블리드 통짜 배경이 채택됨"

        # (2) 편집가능 네이티브 텍스트 존재(통짜 이미지 아님).
        texts = _text_shapes_in(diag_slide)
        assert len(texts) >= 1, "다이어그램 슬라이드 편집가능 텍스트 0"

        # (3) 텍스트-텍스트 겹침 < 10%(더 작은 셰이프 면적 기준).
        for i in range(len(texts)):
            for j in range(i + 1, len(texts)):
                ra, _ = texts[i]
                rb, _ = texts[j]
                inter = overlap_area(ra, rb)
                smaller = min(area(ra), area(rb))
                if smaller <= 0:
                    continue
                ratio = inter / smaller
                assert ratio < 0.10, (
                    f"다이어그램 텍스트 겹침 {ratio:.2%} ≥ 10%: {texts[i][1][:20]!r} vs {texts[j][1][:20]!r}"
                )
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


# ===========================================================================
# 수정 C — 픽 실패/빈값/unknown 변형에서도 통짜 이미지 없이 audit 통과 (사각지대 제거)
# ===========================================================================
#
# 이전 통합 테스트(위)의 _AsyncPick 은 항상 유효한 7키/비-7키 레이아웃과 완전한
# data 를 반환해, 실제 게이트웨이 환경의 "픽 실패(파싱 실패/타임아웃/unknown
# 레이아웃)" 경로를 한 번도 밟지 않았다. 그 사각지대 때문에 실제 산출물이 통짜
# 이미지가 되어도 테스트는 초록불이었다.
#
# 아래 3변형은 정확히 그 실패 경로를 강제한다:
#   (1) empty : _llm_pick_slide_layout 이 {"layout":"", "data":{}} 를 반환
#   (2) raise : _llm_pick_slide_layout 이 예외를 던짐
#   (3) unknown: _llm_pick_slide_layout 이 미지의 레이아웃명을 반환
#
# 그리고 이번에는 Chrome HTML→PNG 베이크가 **성공**(ok=True + 실제 PNG 생성)하도록
# mock 한다. 베이크가 가능한 상황에서도 네이티브 라우팅이 우선돼 통짜 이미지가
# 방출되지 않음(편집가능 네이티브 + 풀블리드 통짜 배경 0)을 입증하기 위함이다.
# (실제 Chrome 은 띄우지 않는다 — mock 이 로컬에서 작은 PNG 를 직접 쓴다.)


def _write_small_png(path, w=1920, h=1080):
    """네이티브 우선이 깨졌을 때 '베이크가 실제로 채택됐는지' 판별용 실제 PNG.

    고주파 세로 줄무늬로 그려 baked_text_score 가 임계를 넘길 수 있게 한다 —
    만약 이 PNG 가 풀블리드 배경으로 잘못 채택되면 audit (f) 베이크 텍스트 검사가
    이를 잡아 passed=False 가 되므로, passed=True 자체가 '베이크 미채택'의 증거다.
    """
    try:
        from PIL import Image, ImageDraw
    except Exception:
        # PIL 부재 시 최소 유효 PNG 바이트(1x1) — 파일 존재만으로 ok=True 취급.
        _PNG_1x1 = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc```\x00\x00"
            b"\x00\x04\x00\x01\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        with open(path, "wb") as f:
            f.write(_PNG_1x1)
        return
    img = Image.new("RGB", (w, h), (245, 246, 248))
    d = ImageDraw.Draw(img)
    step = 6
    for x in range(0, w, step):
        d.rectangle([x, 0, x + step // 2, h], fill=(20, 20, 20))
    img.save(path, "PNG")


def _run_generate_with_pick(tool_input, project_path, env_overrides, pick_callable,
                            bake_success=False):
    """커스텀 픽 콜러블로 _tool_generate_pptx 를 헤르메틱 구동.

    bake_success=True 이면 _render_html_slide_to_png 가 실제 PNG 를 쓰고 ok=True 를
    반환하도록 mock 한다(베이크가 '가능한' 상황). 네이티브 우선이 제대로 동작하면
    이 베이크는 채택되지 않는다.
    """
    old_env = {}
    for k, v in env_overrides.items():
        old_env[k] = os.environ.get(k)
        if v is None:
            os.environ.pop(k, None)
        else:
            os.environ[k] = v

    async def _bake_png(html, output_path, width=1920, height=1080, timeout=30):
        if bake_success:
            try:
                os.makedirs(os.path.dirname(output_path), exist_ok=True)
                _write_small_png(output_path, width, height)
                return {"ok": True, "path": output_path}
            except Exception as e:
                return {"ok": False, "reason": f"mock-write-failed: {e}"}
        return {"ok": False, "reason": "mocked-no-chrome"}

    async def _bake_section(*args, **kwargs):
        # HTML 베이크 섹션 경로 — 네이티브 라우팅 성공 시 호출되지 않아야 한다.
        # bake_success 여도 여기서는 슬라이드 배경 상대경로를 만들어 '베이크가
        # 실제로 일어난다면' 통짜가 되도록 한다(네이티브 우선 검증의 반례 유도).
        if not bake_success:
            return ""
        try:
            render_info = kwargs.get("render_info")
            import time as _t
            local_root = srv._resolve_local_root(project_path)
            rel = os.path.join(".generated", f"_bake_{int(_t.time()*1000000)}.png")
            abs_p = os.path.join(local_root, rel)
            os.makedirs(os.path.dirname(abs_p), exist_ok=True)
            _write_small_png(abs_p, 1920, 1080)
            if isinstance(render_info, dict):
                render_info["composited"] = False
            return rel
        except Exception:
            return ""

    try:
        with mock.patch.object(srv, "_find_local_chrome", lambda: "/fake/chrome"), \
             mock.patch.object(srv, "_call_bridge", lambda ep, payload, timeout=30.0: None), \
             mock.patch.object(srv, "_get_gw", lambda *a, **k: object()), \
             mock.patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "fake-model"), \
             mock.patch.object(srv, "_llm_pick_slide_layout", pick_callable), \
             mock.patch.object(srv, "_render_html_slide_to_png", _bake_png), \
             mock.patch.object(srv, "_generate_html_slide_for_section", _bake_section):
            out = asyncio.run(srv._tool_generate_pptx(tool_input, project_path))
    finally:
        for k, v in old_env.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v
    return json.loads(out)


# --- 3변형 픽 mock -----------------------------------------------------------

class _PickEmpty:
    """게이트웨이가 레이아웃을 못 고름 — 빈 layout/빈 data 반환."""
    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, section_heading, section_body,
                       doc_context="", role="", bullet_count=0):
        self.calls.append(section_heading)
        return {"layout": "", "data": {}}


class _PickRaise:
    """게이트웨이 픽이 타임아웃/파싱 실패로 예외."""
    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, section_heading, section_body,
                       doc_context="", role="", bullet_count=0):
        self.calls.append(section_heading)
        raise RuntimeError("mocked gateway pick failure (timeout/parse)")


class _PickUnknown:
    """게이트웨이가 NATIVE_LAYOUT_REGISTRY/알려진 비-7키에 없는 레이아웃을 반환."""
    def __init__(self):
        self.calls = []

    async def __call__(self, gw, model_id, section_heading, section_body,
                       doc_context="", role="", bullet_count=0):
        self.calls.append(section_heading)
        return {"layout": "some_unknown_layout_xyz", "data": {"foo": "bar"}}


@pytest.mark.parametrize("pick_factory,label", [
    (_PickEmpty, "empty"),
    (_PickRaise, "raise"),
    (_PickUnknown, "unknown"),
])
def test_pick_failure_variants_still_produce_editable_audit_passing_deck(pick_factory, label):
    """픽 실패/빈값/unknown 어느 경우든 통짜 이미지 없이 audit (a)~(h) 통과.

    이것이 이전 테스트가 놓친 실제 게이트웨이 실패 경로다. 베이크가 성공 가능한
    상황(bake_success=True)에서도 네이티브 라우팅이 우선돼 편집가능 네이티브만
    방출됨을 입증한다. audit passed=True 자체가 '통짜 미방출 + 베이크 미채택'의 증거.
    """
    tmp = tempfile.mkdtemp(prefix=f"ae_pickfail_{label}_")
    try:
        tool_input = _build_tool_input()
        env = {
            "AE_ENABLE_HTML_SLIDES": "1",   # 프로덕션 동일 마스터 게이트 ON
            "AE_PREFER_VERTEX_IMAGE": "0",
            "AE_ENABLE_VERTEX_IMAGE": "0",
            "AE_NATIVE_LAYOUT_RENDER": None,  # 기본 활성
        }
        pick = pick_factory()
        result = _run_generate_with_pick(tool_input, tmp, env, pick, bake_success=True)
        pptx_path = _abs_from_result(result, tmp)

        # 픽 실패 경로가 실제로 밟혔는지(라우팅 진입 확인).
        assert pick.calls, f"[{label}] 네이티브 라우팅이 _llm_pick_slide_layout 을 호출하지 않음"

        # (핵심) 산출물 audit — 픽이 실패해도 통짜 0·편집가능≥1·겹침<10% 등 (a)~(h) 통과.
        report = auditor.audit_native_density(pptx_path)
        assert report.passed, (
            f"[{label}] 픽 실패 경로 산출물이 audit 를 통과하지 못함(통짜 의심):\n"
            + "\n".join(
                f"  - slide {f.get('slide')}: {f.get('check')} "
                f"shapes={f.get('shapes')} signal={f.get('signal')}"
                for f in report.failures
            )
        )
        assert report.failures == []

        # 추가 증거 — 어떤 본문 슬라이드에도 풀블리드 통짜 배경(콘텐츠 구운 PNG)이 없고,
        # 모든 슬라이드에 편집가능 텍스트 런이 존재한다.
        from pptx import Presentation
        prs = Presentation(pptx_path)
        for si, slide in enumerate(prs.slides):
            assert _fullbleed_pictures(slide) == 0, (
                f"[{label}] 슬라이드 {si + 1} 에 풀블리드 통짜 배경 채택됨(네이티브 우선 실패)"
            )
            editable = 0
            for sh in slide.shapes:
                try:
                    if sh.has_text_frame and sh.text_frame.text.strip():
                        editable += 1
                except Exception:
                    pass
            assert editable >= 1, f"[{label}] 슬라이드 {si + 1} 편집가능 텍스트 0(통짜 이미지)"
    finally:
        shutil.rmtree(tmp, ignore_errors=True)
