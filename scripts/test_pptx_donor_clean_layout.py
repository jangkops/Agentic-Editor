"""Regression — 템플릿 도너의 빈 콘텐츠 카드/placeholder가 생성 콘텐츠와 겹치지 않는다.

사용자 보고(스크린샷): 템플릿 사용 시 도너 슬라이드의 빈 콘텐츠 카드(둥근 사각형)가 남아
네이티브 다이어그램/불릿과 겹쳐 "오와열이 깨지고 뒤죽박죽"이 됨.

수정: 도너 슬라이드를 배경+로고만 남긴 깨끗한 캔버스로 정리(`_strip_slide_to_background`)한 뒤
제목/다이어그램/불릿을 우리가 정렬 배치. 마무리 단계에서 빈 텍스트 도형 제거
(`_remove_empty_text_shapes`).

Correctness properties:
  P1. 도너 배경(<p:bg>)과 로고(picture)는 보존된다.
  P2. 도너의 빈 콘텐츠 카드/placeholder/샘플 텍스트는 결과물에 남지 않는다.
  P3. 결과 슬라이드에 '빈' 텍스트 도형(보이지 않는 빈 카드/박스)이 남지 않는다.
  P4. 다이어그램형 슬라이드는 편집 가능한 네이티브 도형(텍스트 포함)을 갖는다.

실행: pytest scripts/test_pptx_donor_clean_layout.py -q
"""
from __future__ import annotations

import os
import sys
import io
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")

_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082")


def _make_donor_with_empty_card(path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "SAMPLE_TITLE"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = "SAMPLE_BODY_TEXT"
    # 빈 장식 콘텐츠 카드(텍스트 없는 둥근 사각형) — 겹침 원인
    s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(2), Inches(5), Inches(3))
    # 로고(그림) — 보존되어야 함
    s.shapes.add_picture(io.BytesIO(_PNG), Inches(11.8), Inches(0.2), Inches(1), Inches(1))
    # 배경
    cSld = s._element.find(qn('p:cSld'))
    bg = cSld.makeelement(qn('p:bg'), {})
    bgPr = bg.makeelement(qn('p:bgPr'), {})
    sf = bgPr.makeelement(qn('a:solidFill'), {})
    sf.append(sf.makeelement(qn('a:srgbClr'), {'val': '0B5394'}))
    bgPr.append(sf); bgPr.append(bgPr.makeelement(qn('a:effectLst'), {}))
    bg.append(bgPr); cSld.insert(0, bg)
    prs.save(str(path))


def _slide_stats(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.oxml.ns import qn
    pics = 0
    autoshapes = 0
    empty_text_shapes = 0
    texts = []
    try:
        title_el = slide.shapes.title._element if slide.shapes.title is not None else None
    except Exception:
        title_el = None
    for sh in slide.shapes:
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
            continue
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            autoshapes += 1
        if getattr(sh, "has_text_frame", False):
            t = (sh.text_frame.text or "").strip()
            if t:
                texts.append(t)
            else:
                if title_el is None or sh._element is not title_el:
                    empty_text_shapes += 1
    has_bg = slide._element.find(qn('p:cSld')).find(qn('p:bg')) is not None
    return {"pics": pics, "autoshapes": autoshapes,
            "empty_text_shapes": empty_text_shapes, "text": "\n".join(texts), "has_bg": has_bg}


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_diagram_slide_on_donor_is_clean(tmp_path):
    """P1~P4 — 다이어그램 슬라이드: 배경/로고 보존, 빈 카드 제거, 네이티브 도형 존재."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    tpl = tmp_path / "donor.pptx"
    _make_donor_with_empty_card(tpl)

    slides = [{"title": "데이터 흐름", "bullets": ["입력 -> 처리 -> 출력"],
               "nativeDiagram": {"type": "flow", "content": "입력 -> 처리 -> 출력"}}]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "새 덱", "slides": slides, "templatePath": str(tpl)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        from pptx import Presentation
        prs = Presentation(res["absPath"])
        content = list(prs.slides)[1]  # cover=0, content=1
        st = _slide_stats(content)
        assert st["has_bg"], "배경 손실"
        assert st["pics"] >= 1, "로고(그림) 손실"
        assert st["autoshapes"] > 0, "네이티브 다이어그램 도형 없음"
        assert st["empty_text_shapes"] == 0, f"빈 텍스트 도형 잔존: {st['empty_text_shapes']}"
        assert "SAMPLE" not in st["text"], "도너 샘플 텍스트 잔존"
    finally:
        _cleanup(res)


def test_text_slide_on_donor_is_clean(tmp_path):
    """불릿(텍스트) 슬라이드: 빈 카드 제거, 불릿은 깨끗한 텍스트박스에 배치."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    tpl = tmp_path / "donor2.pptx"
    _make_donor_with_empty_card(tpl)

    slides = [{"title": "개요", "bullets": ["항목 가", "항목 나"]}]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "새 덱", "slides": slides, "templatePath": str(tpl)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        from pptx import Presentation
        prs = Presentation(res["absPath"])
        content = list(prs.slides)[1]
        st = _slide_stats(content)
        assert st["has_bg"]
        assert st["pics"] >= 1
        assert st["empty_text_shapes"] == 0, f"빈 텍스트 도형 잔존: {st['empty_text_shapes']}"
        assert "SAMPLE" not in st["text"]
        assert "항목 가" in st["text"] and "개요" in st["text"]
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
