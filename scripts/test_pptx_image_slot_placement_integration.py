"""통합 검증 테스트 — spec: pptx-image-slot-placement-fix (bugfix), Task 5.

실제 ``ai_engine/server.py:_tool_generate_pptx`` 합성 경로를 **헤르메틱 목**(네트워크 0)
으로 끝까지 구동해, 본 스펙의 결함 3종이 **덱 레벨**에서 해소됐는지와 손실-0 보존을
audit 도구로 검사한다. 자매 스펙 ``test_pptx_quality_vertex_images_integration.py`` 의
헤르메틱 목 패턴(``_FakeVertexClient``/``_render_html_png_fake``/``_img_gen_disabled`` +
``patch.object`` 로 ``_call_bridge``/``_get_gw``/``_render_html_slide_to_png``/
``_generate_html_slide_for_section`` 목, Vertex 는 ``get_vertex_image_client`` 스텁 +
``generate`` 목)을 재사용한다. 측정은 기존 audit 도구
(``audit_pptx_zorder_break.py``/``audit_pptx_media_classify.py``)의 공개 판정 함수를
재사용해 탐색·검증·통합이 동일 기준(풀블리드 0.92 비율, off-slide EPS=0.05)으로
측정되게 한다.

  - **D1 통합** (design Property 1, Req 2.1) — ``coverBackground`` + HTML 표지를 동시
    활성한 덱을 합성 → 생성 덱의 **각 슬라이드 풀블리드 PICTURE ≤ 1**
    (``audit_pptx_zorder_break._fullbleed`` / ``audit_pptx_media_classify._fb`` 로 검증).
    수정(server.py 4836 ``and not _cover_bg_embedded`` + 본문 ``_fb_embedded`` 가드)이
    중복 풀블리드를 차단한다.

  - **D2 통합** (design Property 2, Req 2.2) — 대형(3840×2160) 이미지가 카드 아이콘 칩
    (≤0.5in) 슬롯으로 흘러드는 입력을 합성 → 생성 덱에 **소형 슬롯의 대형 이미지 0건**
    (``audit_pptx_media_classify`` 의 슬롯-이미지 정합 분류 = PICTURE 픽셀 크기 ×
    배치 슬롯 크기로 검증). 수정(``native_diagram_pptx`` 아이콘 칩 ``slot_image_fits``
    검사 → 대형이면 PNG 스킵·글리프 폴백)이 오배정을 막는다.

  - **D3 통합** (design Property 3, Req 2.3) — 큰 부분 이미지(``imageFile``)를 가진 덱을
    합성 → ``audit_pptx_zorder_break`` 의 **off-slide(슬라이드 밖) 검출 0건**. 수정
    (server.py ~5475 ``clamp_into_bounds`` 적용)이 모든 부분 이미지를 경계 안으로
    클램프한다.

  - **손실-0 (P5)** (design Property 5, Req 3.2) — Vertex 이미지 생성 입력 → 생성된
    **모든 Vertex 이미지가 ppt/media 에 임베드**(unused == 0). 덱 레벨 보조 단언으로
    ``_select_render_plan`` 의 비주얼 슬라이드 결정 ``vertex_slot != "none"`` 도 확인.

전부 헤르메틱 — 게이트웨이(``_get_gw``/``_call_bridge`` 스파이), Vertex
(``get_vertex_image_client`` 스텁 + ``generate`` 목), HTML→PNG 렌더 모두 목 처리.
``heredoc``/``stdin`` 미사용 — 파일로 작성해 pytest 로 실행한다.

실행(헤르메틱):
  ./venv/bin/python -m pytest scripts/test_pptx_image_slot_placement_integration.py -p no:cacheprovider -q

_Requirements: 2.1, 2.2, 2.3, 3.2_
"""
from __future__ import annotations

import io
import os
import sys
import json
import base64
import asyncio
import zipfile
import tempfile
from contextlib import ExitStack, redirect_stdout
from unittest.mock import patch

# repo root + scripts dir 를 import 경로에 추가.
_HERE = os.path.dirname(__file__)
sys.path.insert(0, os.path.join(_HERE, ".."))
sys.path.insert(0, _HERE)

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402

# audit 도구의 공개 판정 함수 재사용(동일 기준 측정).
import audit_pptx_zorder_break as azb  # noqa: E402
import audit_pptx_media_classify as acm  # noqa: E402

# --------------------------------------------------------------------------
# 임계 상수 — design §0 / Bug Condition 보조 술어와 동일.
# --------------------------------------------------------------------------
SW, SH = 13.333, 7.5
LARGE_PX = 1024
SMALL_SLOT_IN = 0.5
EPS = 0.05


# --------------------------------------------------------------------------
# 헤르메틱 fakes (자매 스펙 test_pptx_quality_vertex_images_integration.py 패턴)
# --------------------------------------------------------------------------
def _make_png(path, w, h, tag=0):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    Image.new("RGB", (w, h),
              (tag % 256, (tag * 7) % 256, (tag * 31) % 256)).save(path, "PNG")
    return path


def _make_png_bytes(tag: int) -> bytes:
    """각 생성 이미지가 바이트-고유하도록 unique PNG 생성."""
    img = Image.new("RGB", (40 + (tag % 11), 30 + (tag % 7)),
                    (tag % 256, (tag * 7) % 256, (tag * 31) % 256))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class _FakeVertexClient:
    """VertexImageClient 대역 — 항상 'enabled', 매 호출 바이트-고유 PNG 반환.

    생성한 raw 바이트를 기록해 손실-0(모든 생성 이미지 임베드)을 바이트 레벨로 검증."""

    def __init__(self) -> None:
        self.enabled = True
        self.calls = 0
        self.generated_raw: list[bytes] = []

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.calls += 1
        raw = _make_png_bytes(2000 + self.calls)
        self.generated_raw.append(raw)
        return {"images": [base64.b64encode(raw).decode("ascii")]}


class _DisabledVertexClient:
    """Vertex 비활성 스텁 — 네트워크 0, 이미지 생성 없음(D1/D2/D3 시나리오)."""
    enabled = False

    async def generate(self, *a, **k):
        return {"images": []}


async def _img_gen_disabled(*_a, **_k):
    """_tool_generate_image 대역 — 경로 반환 안 함(네트워크 0)."""
    return json.dumps({"error": "test-disabled"})


def _make_html_png_fake(tag):
    async def _fake(html, output_path, width=1920, height=1080, timeout=30, **_k):
        _make_png(output_path, 64, 36, tag)
        return {"ok": True}
    return _fake


def _make_section_fake(project_path: str, tag_base: int):
    """섹션-HTML 렌더 대역 — 풀블리드 슬라이드배경 PNG 를 .generated 에 기록."""
    state = {"n": 0}

    async def _section_fake(gw, model_id, heading, body, doc_context,
                            proj_path, style_profile=None,
                            hero_image="", render_info=None):
        state["n"] += 1
        gen = os.path.join(project_path, ".generated")
        os.makedirs(gen, exist_ok=True)
        name = f"slide-html-{state['n']}.png"
        out = os.path.join(gen, name)
        _make_png(out, 1920, 1080, tag_base + state["n"])
        if isinstance(render_info, dict):
            render_info["layout"] = "feature_grid"
            render_info["composited"] = False
        return f".generated/{name}"

    return _section_fake


# --------------------------------------------------------------------------
# audit 기반 검사 helpers — audit 도구의 공개 함수 재사용
# --------------------------------------------------------------------------
def _pic_pixels(sh):
    """PICTURE 의 원본 픽셀 크기(없으면 (0,0))."""
    try:
        with Image.open(io.BytesIO(sh.image.blob)) as im:
            return im.size
    except Exception:
        return (0, 0)


def _slide_pictures(slide):
    """슬라이드의 PICTURE 목록 → [(rect, (px_w, px_h)), ...] (z-order 순).

    rect 측정은 ``audit_pptx_zorder_break._rect`` 재사용."""
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        out.append((azb._rect(sh), _pic_pixels(sh)))
    return out


def _fullbleed_counts(pptx_abs):
    """슬라이드별 풀블리드 PICTURE 개수(``audit`` 의 ``_fullbleed`` 재사용)."""
    counts = []
    for slide in Presentation(pptx_abs).slides:
        counts.append(sum(1 for r, _px in _slide_pictures(slide)
                          if azb._fullbleed(r)))
    return counts


def _fullbleed_counts_media_classify(pptx_abs):
    """교차검증 — ``audit_pptx_media_classify._fb`` 로 슬라이드별 풀블리드 개수."""
    counts = []
    for slide in Presentation(pptx_abs).slides:
        n = 0
        for sh in slide.shapes:
            try:
                if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue
            if acm._fb(acm._rect(sh)):
                n += 1
        counts.append(n)
    return counts


def _large_in_small_slot(pptx_abs):
    """소형 슬롯(≤0.5in)에 대형 이미지(≥1024px)가 배정된 PICTURE 목록.

    슬롯-이미지 정합 분류: 배치 rect 는 ``audit_pptx_media_classify._rect`` 로,
    이미지 픽셀 크기는 blob 디코드로 측정한다."""
    offenders = []
    for idx, slide in enumerate(Presentation(pptx_abs).slides):
        for sh in slide.shapes:
            try:
                if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue
            r = acm._rect(sh)
            pw, ph = _pic_pixels(sh)
            is_large = (pw >= LARGE_PX or ph >= LARGE_PX)
            is_small = (r is not None and r[2] <= SMALL_SLOT_IN
                        and r[3] <= SMALL_SLOT_IN)
            if is_large and is_small:
                offenders.append((idx, r, (pw, ph)))
    return offenders


def _offslide_count(pptx_abs):
    """슬라이드 밖 PICTURE 개수 — ``audit_pptx_zorder_break`` 의 off-slide 판정과
    동일 정의(``r[0]<-0.05 ∨ r[1]<-0.05 ∨ r[0]+r[2]>SW+0.05 ∨ r[1]+r[3]>SH+0.05``)."""
    n = 0
    for slide in Presentation(pptx_abs).slides:
        for r, _px in _slide_pictures(slide):
            if r is None:
                continue
            if (r[0] < -EPS or r[1] < -EPS
                    or r[0] + r[2] > SW + EPS or r[1] + r[3] > SH + EPS):
                n += 1
    return n


def _run_zorder_audit(pptx_abs):
    """``audit_pptx_zorder_break.audit`` 를 실제 호출(narration) — stdout 캡처."""
    buf = io.StringIO()
    with redirect_stdout(buf):
        azb.audit(pptx_abs)
    return buf.getvalue()


def _pptx_media_bytes(pptx_abs):
    with zipfile.ZipFile(pptx_abs) as z:
        return [z.read(n) for n in z.namelist() if n.startswith("ppt/media/")]


def _vertex_usage(fake, pptx_abs):
    """(generated, embedded, unused) — 생성 Vertex 이미지의 임베드 현황."""
    media = _pptx_media_bytes(pptx_abs)
    embedded = sum(1 for raw in fake.generated_raw if raw in media)
    generated = len(fake.generated_raw)
    return generated, embedded, generated - embedded


def _big_icon_factory(big_path):
    def _fake_get_icon_png(name, color_hex="FFFFFF", px=240):
        # 정상은 소형 아이콘 PNG — 여기서는 대형 4K 이미지가 칩 슬롯으로 흘러드는
        # 상황을 재현해 수정 가드(slot_image_fits)가 막는지 확인.
        return big_path
    return _fake_get_icon_png


# ==========================================================================
# D1 통합 — coverBackground + HTML 표지 동시 → 각 슬라이드 풀블리드 ≤ 1
# ==========================================================================
def test_d1_integration_each_slide_fullbleed_le_1():
    """슬라이드 8·9 유사: 표지 배경 + HTML 표지 동시 활성 덱 합성 →
    각 슬라이드 풀블리드 PICTURE ≤ 1 (audit 두 도구로 교차검증)."""
    proj = tempfile.mkdtemp()
    cover_bg = _make_png(os.path.join(proj, "cover-bg.png"), 1920, 1080, tag=11)

    env = {
        "AE_ENABLE_HTML_SLIDES": "1",       # HTML 표지/본문 풀블리드 경로 활성
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_ENABLE_VERTEX_BG": "0",         # Vertex 공유 배경 OFF(coverBackground 보존)
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {
        "title": "표지 풀블리드 중복 통합 검증",
        "coverBackground": cover_bg,        # 1번째 풀블리드 후보
        "slides": [
            {"title": "본문 1", "bullets": ["요점 A", "요점 B"]},
            {"title": "본문 2", "bullets": ["요점 C", "요점 D"]},
        ],
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: _DisabledVertexClient()), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", lambda *a, **k: object()), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_render_html_slide_to_png", _make_html_png_fake(21)), \
            patch.object(srv, "_generate_html_slide_for_section", _make_section_fake(proj, 600)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    pptx = result["absPath"]

    fb_zorder = _fullbleed_counts(pptx)
    fb_mclass = _fullbleed_counts_media_classify(pptx)
    print(f"[D1] 풀블리드 count (zorder audit)        = {fb_zorder}")
    print(f"[D1] 풀블리드 count (media_classify audit) = {fb_mclass}")

    worst = max(fb_zorder) if fb_zorder else 0
    assert worst <= 1, (
        f"D1 위반: 한 슬라이드에 풀블리드 배경 {worst}장(기대 ≤1). "
        f"슬라이드별 count={fb_zorder}")
    # 두 audit 도구가 동일 기준으로 동일 결과를 줘야 한다(감사↔코드 일치).
    assert fb_zorder == fb_mclass, (
        f"audit 도구 간 풀블리드 판정 불일치: zorder={fb_zorder} mclass={fb_mclass}")


# ==========================================================================
# D2 통합 — 대형 이미지 입력 → 소형 슬롯에 대형 이미지 0건
# ==========================================================================
def test_d2_integration_no_large_image_in_small_slot():
    """슬라이드 8·9 유사: 대형(3840×2160) 이미지가 카드 아이콘 칩(≤0.5in) 슬롯으로
    흘러드는 입력 합성 → 소형 슬롯의 대형 이미지 0건(media_classify 정합 분류)."""
    proj = tempfile.mkdtemp()
    big = _make_png(os.path.join(proj, "big-4k.png"), 3840, 2160, tag=33)

    env = {
        "AE_ENABLE_HTML_SLIDES": "0",       # HTML OFF → nativeDiagram 카드 경로 직행
        "AE_DISABLE_HTML_SLIDES": "1",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_DISABLE_RICH_ICONS": "0",       # 아이콘 칩 PNG 경로 활성
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    cards = "\n".join(f"기능{i}: 핵심 설명 {i}" for i in range(1, 7))  # 6 카드 → cols=3
    deck = {
        "title": "대형이미지 소형슬롯 통합 검증",
        "slides": [
            {"title": "핵심 기능", "nativeDiagram": {"type": "cards", "content": cards}},
        ],
    }

    # native_diagram_pptx 는 함수 내부에서 icon_assets.get_icon_png 를 호출 — 양쪽 패치.
    patchers = []
    try:
        import ai_engine.icon_assets as _ia_pkg
        patchers.append(patch.object(_ia_pkg, "get_icon_png", _big_icon_factory(big)))
    except Exception:
        pass
    try:
        import icon_assets as _ia_top  # noqa
        patchers.append(patch.object(_ia_top, "get_icon_png", _big_icon_factory(big)))
    except Exception:
        pass
    assert patchers, "icon_assets 모듈 패치 실패"

    with ExitStack() as stack:
        for p in patchers:
            stack.enter_context(p)
        stack.enter_context(patch.dict(os.environ, env, clear=False))
        stack.enter_context(patch.object(vim, "get_vertex_image_client", lambda **_k: _DisabledVertexClient()))
        stack.enter_context(patch.object(srv, "_call_bridge", lambda *a, **k: None))
        stack.enter_context(patch.object(srv, "_find_local_chrome", lambda: ""))
        stack.enter_context(patch.object(srv, "_tool_generate_image", _img_gen_disabled))
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    pptx = result["absPath"]

    offenders = _large_in_small_slot(pptx)
    print(f"[D2] 소형 슬롯의 대형 이미지 = {offenders}")
    assert not offenders, (
        f"D2 위반: 대형 이미지가 소형 장식 슬롯에 배정됨 — {offenders} "
        f"(슬롯-이미지 크기 정합 가드 부재).")


# ==========================================================================
# D3 통합 — 큰 부분 이미지 입력 → off-slide 검출 0건
# ==========================================================================
def test_d3_integration_offslide_detection_zero():
    """슬라이드 1 유사: 큰 부분 이미지(imageFile)를 가진 덱 합성 →
    audit_pptx_zorder_break 의 off-slide 검출 0건(clamp_into_bounds 적용)."""
    proj = tempfile.mkdtemp()
    # region 보다 큰 세로/가로 이미지 — 부분-이미지 배치 경로(img_path) 구동.
    illo_tall = _make_png(os.path.join(proj, "illo-tall.png"), 1200, 2400, tag=44)
    illo_wide = _make_png(os.path.join(proj, "illo-wide.png"), 3000, 900, tag=45)

    env = {
        "AE_ENABLE_HTML_SLIDES": "0",        # HTML OFF → 풀블리드 bg_path 미사용
        "AE_DISABLE_HTML_SLIDES": "1",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # nativeDiagram 추론 skip → imageFile 보존
        "AE_DISABLE_NATIVE_DIAGRAM": "1",
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {
        "title": "부분 이미지 경계 통합 검증",
        "slides": [
            {"title": "세로 일러스트", "bullets": ["설명 A", "설명 B"],
             "imageFile": illo_tall},
            {"title": "가로 일러스트", "imageFile": illo_wide},
        ],
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: _DisabledVertexClient()), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    pptx = result["absPath"]

    # audit 도구 실제 호출(narration) — off-slide 요약을 stdout 으로 확인.
    audit_out = _run_zorder_audit(pptx)
    print(audit_out)

    offslide = _offslide_count(pptx)
    print(f"[D3] off-slide PICTURE count = {offslide}")
    assert "슬라이드 밖 이미지     : 0" in audit_out, (
        f"D3 위반: audit off-slide 요약이 0이 아님\n{audit_out}")
    assert offslide == 0, (
        f"D3 위반: 부분 이미지가 슬라이드 경계 밖에 배치됨 — off-slide={offslide}건. "
        f"clamp_into_bounds 경계 가드가 적용되지 않았다.")


# ==========================================================================
# 손실-0 (P5) — Vertex 이미지 생성 입력 → 모든 생성 이미지 임베드(unused == 0)
# ==========================================================================
def test_lossless_p5_all_vertex_images_embedded():
    """비주얼/콘텐츠 슬라이드에 Vertex 이미지 생성 → 모든 생성 이미지 ppt/media 임베드
    (unused == 0). 덱 레벨 보조 단언: _select_render_plan 비주얼 결정 vertex_slot != none."""
    fake = _FakeVertexClient()
    proj = tempfile.mkdtemp()

    env = {
        "AE_ENABLE_HTML_SLIDES": "0",        # HTML off → Vertex 이미지 바이트로 임베드
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex 사전생성 ON
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # 결정론: 게이트웨이 구조화 없음
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {
        "title": "손실-0 통합 검증",
        "slides": [
            {"title": "회사 소개", "bullets": ["신뢰를 최우선으로"],
             "imagePrompt": "a high quality professional photograph of a modern "
                            "corporate office, natural light, wide angle"},
            {"title": "비전", "bullets": ["미래를 향한 도전"],
             "imagePrompt": "an inspiring wide landscape photo, golden hour"},
        ],
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    pptx = result["absPath"]

    generated, embedded, unused = _vertex_usage(fake, pptx)
    print(f"[P5] Vertex generated={generated} embedded={embedded} unused={unused}")
    assert generated >= 1, (
        f"비주얼 슬라이드에 Vertex 이미지가 생성돼야 함 (generated={generated})")
    assert unused == 0, (
        f"손실-0 위반: 생성된 Vertex 이미지가 폐기됨 "
        f"(generated={generated}, embedded={embedded}, unused={unused})")

    # 덱 레벨 보조 단언 — 순수 결정 함수가 비주얼 슬라이드를 폐기(none)하지 않는다.
    plan = srv._select_render_plan(
        has_vertex_image=True, has_native_diagram=False,
        has_image_file=False, has_slide_bg=False,
        role="visual", html_enabled=False,
    )
    print(f"[P5] _select_render_plan(visual) = {plan}")
    assert plan["vertex_slot"] != "none", (
        f"손실-0 위반: 비주얼 슬라이드의 Vertex 이미지가 'none' 슬롯으로 폐기됨 — {plan}")


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
