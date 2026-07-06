"""Regression — Genspark급 HTML 디자인 슬라이드가 직접 생성 경로에서 풀블리드로 적용된다.

사용자 요구: "고퀄리티로 젠스파크 이미지 슬라이드 급" 결과물.

수정: `_tool_generate_pptx`가 무템플릿 + Electron 브리지(render-html-to-png) 가용 시
각 슬라이드(표지 포함)를 HTML 디자인 레이아웃으로 렌더해 풀블리드 배경으로 사용한다.
브리지가 없으면(헤드리스/테스트) 네이티브 도형 경로로 자동 폴백한다(회귀 없음).

Correctness properties:
  P1. 브리지 가용 시: 모든 슬라이드가 풀블리드 HTML PNG(슬라이드 전체를 덮는 Picture)를 갖는다.
  P2. HTML 슬라이드가 적용되면 네이티브 다이어그램 추론은 일어나지 않는다.
  P3. 브리지 미가용 시: HTML 호출 없이 기존 네이티브 경로로 폴백한다.

실행: pytest scripts/test_pptx_html_slides.py -q
"""
from __future__ import annotations

import os
import sys
import io
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")

_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082")


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_html_slides_used_when_bridge_available(tmp_path, monkeypatch):
    """P1+P2 — 브리지 가용 시 모든 슬라이드가 풀블리드 HTML PNG로 채워진다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "1")  # 이 파일은 HTML 슬라이드 기능을 검증(기본 OFF를 opt-in)
    gen = tmp_path / ".generated"
    gen.mkdir(parents=True, exist_ok=True)

    # 브리지 가용으로 위장
    monkeypatch.setattr(server, "_call_bridge", lambda ep, payload, timeout=30.0: {"remote": False})
    monkeypatch.setattr(server, "_get_gw", lambda *a, **k: object())
    monkeypatch.setattr(server, "_specialized_model_for_task",
                        lambda *a, **k: "us.anthropic.claude-sonnet-4-6")

    calls = {"n": 0}

    async def _fake_html(gw, model, heading, body, ctx, project_path, style_profile=None):
        calls["n"] += 1
        fn = gen / f"html-slide-{calls['n']}.png"
        fn.write_bytes(_PNG)
        return f".generated/{fn.name}"

    monkeypatch.setattr(server, "_generate_html_slide_for_section", _fake_html)
    # 네이티브 다이어그램이 절대 호출되지 않아야 함(HTML 우선)
    import native_diagram_pptx as nd
    monkeypatch.setattr(nd, "build_native_diagram",
                        lambda *a, **k: (_ for _ in ()).throw(AssertionError("native diagram should not run")))

    slides = [
        {"title": "아키텍처", "bullets": ["프론트 -> 백엔드 -> DB"]},
        {"title": "데이터 흐름", "bullets": ["입력 -> 처리 -> 출력"]},
    ]
    out = asyncio.run(server._tool_generate_pptx({"title": "프로젝트 개요", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        # cover + 2 content = 3, 각 슬라이드에 풀블리드 그림이 있어야 함
        assert len(prs.slides._sldIdLst) == 3
        slide_w, slide_h = prs.slide_width, prs.slide_height
        for i, sl in enumerate(prs.slides):
            pics = [sh for sh in sl.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            assert pics, f"슬라이드 {i}에 HTML 풀블리드 그림 없음"
            # 풀블리드: 슬라이드 크기의 ~95% 이상 덮는 그림이 하나 이상
            full = [p for p in pics if p.width >= slide_w * 0.95 and p.height >= slide_h * 0.95]
            assert full, f"슬라이드 {i} 그림이 풀블리드가 아님"
        # 표지(1) + 콘텐츠(2) = 3회 HTML 렌더 호출
        assert calls["n"] == 3, f"HTML 렌더 호출 수={calls['n']} (기대 3)"
    finally:
        _cleanup(res)


def test_falls_back_to_native_when_bridge_unavailable(tmp_path, monkeypatch):
    """P3 — 브리지 미가용 시 HTML 호출 없이 네이티브 경로로 폴백."""
    from pptx import Presentation
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "1")  # 이 파일은 HTML 슬라이드 기능을 검증(기본 OFF를 opt-in)

    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)  # 브리지 없음
    called = {"n": 0}

    async def _fake_html(*a, **k):
        called["n"] += 1
        return ""

    monkeypatch.setattr(server, "_generate_html_slide_for_section", _fake_html)

    slides = [{"title": "개요", "bullets": ["항목 1", "항목 2"]}]
    out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        assert called["n"] == 0, "브리지 없는데 HTML 렌더가 호출됨"
        prs = Presentation(res["absPath"])
        assert len(prs.slides._sldIdLst) == 2
    finally:
        _cleanup(res)


def test_template_uses_html_with_style_profile(tmp_path, monkeypatch):
    """템플릿 사용 시에도 HTML 고품질 렌더를 사용한다(젠스파크급 레이아웃 + 템플릿 색/폰트).

    근거(사용자 결정): used_template일 때 HTML을 끄면 항상 휑한 네이티브 도너로 빠져
    품질이 급락. 진단 로그(used_template=True → _html_enabled=False)로 근본 원인 확정.
    따라서 템플릿이어도 HTML 렌더를 쓰되 Style_Profile(색/폰트)을 HTML 디자인 토큰에
    주입한다(_generate_html_slide_for_section의 style_profile 인자).

    Correctness properties:
      T1. 템플릿 + 브리지 가용 시 _generate_html_slide_for_section이 호출된다(HTML 사용).
      T2. style_profile이 HTML 렌더 호출에 전달된다(템플릿 색 반영).
      T3. 출력 슬라이드가 풀블리드 HTML PNG를 갖는다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches
    from pptx.oxml.ns import qn

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "1")
    gen = tmp_path / ".generated"; gen.mkdir(parents=True, exist_ok=True)

    tpl = Presentation(); tpl.slide_width = Inches(13.333); tpl.slide_height = Inches(7.5)
    for i in range(3):
        s = tpl.slides.add_slide(tpl.slide_layouts[1])
        s.shapes.title.text = f"TPL_SAMPLE_{i}"
    tpl_path = tmp_path / "tpl.pptx"
    tpl.save(str(tpl_path))

    monkeypatch.setattr(server, "_call_bridge", lambda ep, payload, timeout=30.0: {"remote": False})
    monkeypatch.setattr(server, "_get_gw", lambda *a, **k: object())
    monkeypatch.setattr(server, "_specialized_model_for_task", lambda *a, **k: "us.anthropic.claude-sonnet-4-6")

    seen = {"n": 0, "profiles": []}

    async def _fake_html(gw, model, heading, body, ctx, project_path, style_profile=None):
        seen["n"] += 1
        seen["profiles"].append(style_profile)
        fn = gen / f"h-{seen['n']}.png"
        fn.write_bytes(_PNG)
        return f".generated/{fn.name}"

    monkeypatch.setattr(server, "_generate_html_slide_for_section", _fake_html)

    style_profile = {"primaryColor": "#0B5394", "textColor": "#1A1A1A",
                     "headingFont": "Pretendard", "bodyFont": "Pretendard"}
    slides = [{"title": "흐름", "bullets": ["A -> B"]}, {"title": "구조", "bullets": ["X"]}]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "덱", "slides": slides,
         "templatePath": str(tpl_path), "styleProfile": style_profile}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        # T1 — 템플릿이어도 HTML 렌더가 호출된다 (표지 + 콘텐츠)
        assert seen["n"] >= 1, "템플릿인데 HTML 렌더가 호출되지 않음(휑한 네이티브로 빠짐)"
        # T2 — style_profile이 전달된다(템플릿 색 반영)
        assert any(p == style_profile for p in seen["profiles"]), \
            f"HTML 렌더에 style_profile 미전달: {seen['profiles']}"
        # T3 — 풀블리드 HTML PNG가 슬라이드에 임베드된다
        prs = Presentation(res["absPath"])
        slide_w, slide_h = prs.slide_width, prs.slide_height
        full = sum(1 for sl in prs.slides for sh in sl.shapes
                   if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
                   and sh.width >= slide_w * 0.95 and sh.height >= slide_h * 0.95)
        assert full >= 1, "템플릿+HTML인데 풀블리드 PNG가 없음"
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
