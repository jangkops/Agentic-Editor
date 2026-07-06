"""Regression — 템플릿 사용 시 본문 내용이 누락되지 않는다(부분 누락 버그 수정).

사용자 보고: "템플릿 사용하면 부분만 가져오고 부분은 누락된다."
근본 원인: _strip_slide_to_background가 도너의 본문 placeholder까지 전부 제거 →
텍스트 슬라이드의 불릿을 채울 그릇이 없어 내용이 사라졌다.

수정: placeholder는 '그릇'으로 보존(샘플 텍스트만 비움), 비-placeholder 장식 텍스트박스만 제거.

Correctness properties:
  P1. 템플릿 텍스트 슬라이드의 새 불릿이 결과물에 존재한다(누락 없음).
  P2. 템플릿의 샘플 텍스트(마커)는 남지 않는다.
  P3. 도너의 장식 비-placeholder 텍스트박스(샘플 카드)는 제거된다.

실행: pytest scripts/test_pptx_template_body_preserved.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_template_text_slide_keeps_bullets(tmp_path):
    """P1+P2+P3 — 템플릿 텍스트 슬라이드 불릿 보존 + 샘플 제거."""
    from pptx import Presentation
    from pptx.util import Inches

    # 템플릿: 제목+본문 placeholder(샘플 텍스트) + 장식 텍스트박스(비-placeholder)
    tpl = Presentation()
    tpl.slide_width = Inches(13.333); tpl.slide_height = Inches(7.5)
    for i in range(2):
        s = tpl.slides.add_slide(tpl.slide_layouts[1])  # Title and Content
        s.shapes.title.text = f"TPL_TITLE_{i}"
        if len(s.placeholders) > 1:
            s.placeholders[1].text = f"TPL_SAMPLE_BODY_{i}"
        # 장식 비-placeholder 텍스트박스(샘플 카드)
        tb = s.shapes.add_textbox(Inches(9), Inches(0.4), Inches(3), Inches(0.6))
        tb.text_frame.text = f"DECOR_CARD_{i}"
    tpl_path = tmp_path / "tpl.pptx"
    tpl.save(str(tpl_path))

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    # 다이어그램/이미지 경로를 타지 않도록 끔 → 순수 텍스트 슬라이드 검증
    os.environ["AE_DISABLE_NATIVE_DIAGRAM"] = "1"
    slides = [
        {"title": "신규 섹션 A", "bullets": ["핵심 항목 가나다", "두번째 항목 라마바"]},
        {"title": "신규 섹션 B", "bullets": ["세번째 항목 사아자"]},
    ]
    try:
        out = asyncio.run(server._tool_generate_pptx(
            {"title": "새 덱", "slides": slides, "templatePath": str(tpl_path)}, ""))
        res = json.loads(out)
        try:
            assert "error" not in res, res
            prs = Presentation(res["absPath"])
            alltext = "\n".join(
                sh.text_frame.text
                for sl in prs.slides for sh in sl.shapes
                if sh.has_text_frame)
            # P1 — 새 불릿 보존(누락 없음)
            for b in ("핵심 항목 가나다", "두번째 항목 라마바", "세번째 항목 사아자"):
                assert b in alltext, f"본문 불릿 누락: {b}"
            # P2 — 템플릿 샘플 본문 제거
            assert "TPL_SAMPLE_BODY" not in alltext, "템플릿 샘플 본문이 잔존"
            # P3 — 장식 카드 제거
            assert "DECOR_CARD" not in alltext, "장식 텍스트박스가 잔존(클러터)"
        finally:
            _cleanup(res)
    finally:
        os.environ.pop("AE_DISABLE_NATIVE_DIAGRAM", None)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
