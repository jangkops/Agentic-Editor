"""Preservation + gateway-constraint property tests — spec: pptx-quality-vertex-images
(bugfix), Tasks 2 and 6.

PROPERTY 2 — Preservation (구조형/무관 입력의 기존 동작 보존).
PROPERTY 5 — Gateway Constraint (이미지 외 호출은 Vertex 미사용) — appended in Task 6;
see the PROPERTY 5 section near the bottom of this file (PROP5-1 / PROP5-2 / PROP5-3).

These property-based tests pin down the behaviours the fix (tasks 3.x) MUST NOT
change. Following the **observation-first** methodology, every assertion below
was derived by first OBSERVING what the *current* (unfixed) ``ai_engine/server.py``
``_tool_generate_pptx`` decision seam actually does on NON-bug
(``isBugCondition == False``) inputs, then asserting exactly that. They therefore
PASS on the unfixed code (establishing the baseline) and must keep PASSing after
the fix (regression guard).

NON-bug input families covered (each precondition-checked against the design
``isBugCondition`` mirror so the test stays honest):

  PRES-1  구조형 보존 (Req 3.1)
      A real structural slide (flow / tree / architecture) renders as EDITABLE
      native shapes: Vertex pre-gen is skipped (``_gen_vertex_slide`` returns ""
      for a structural kind), NO raster image is embedded (``ppt/media`` empty),
      and the bullet labels appear as native shape text (not rasterised).

  PRES-2  HTML 풀블리드 보존 (Req 3.2)
      With HTML full-bleed enabled and Vertex DISABLED (so ``isBugCondition`` is
      false — no gateSuppressed), a content slide takes the ``slideBackground``
      path: the section HTML→PNG capture is embedded as a (0,0) full-bleed
      PICTURE covering the whole 13.333×7.5in slide. The path is unchanged.

  PRES-3  Vertex 비활성/실패 폴백 보존 (Req 3.3)
      When Vertex is disabled OR its ``generate`` fails, ``_vertex_pre`` stays
      empty and a non-structural content slide falls back to the native text
      path: NO raster media is embedded and the body bullet text is preserved in
      a placeholder (media-output-quality regression prevention).

  PRES-4  템플릿 상속 보존 (Req 3.5)
      When ``styleProfile`` is supplied, the per-section HTML renderer is invoked
      WITH that exact profile (HTML design-token inheritance), and
      ``_build_palette(styleProfile)`` yields a stable native palette
      (native-palette inheritance). Neither path is altered.

  PRES-5  명시 우선순위 보존 (caller imageFile / slideBackground)
      A slide whose caller already specified ``imageFile`` (or ``slideBackground``)
      keeps that exact image: Vertex pre-gen skips it and the caller's bytes are
      the ones embedded — existing precedence is unchanged.

Everything is hermetic — no network. The Bedrock gateway, the Vertex client
(``get_vertex_image_client`` / ``generate``), the HTML→PNG renderer and
``_tool_generate_image`` are all mocked. We open the produced ``.pptx`` as a zip
and inspect ``ppt/media/*`` bytes + shape geometry to assert behaviour.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_quality_vertex_images_preservation_pbt.py -p no:cacheprovider -q

_Requirements: 3.1, 3.2, 3.3, 3.5_
"""
from __future__ import annotations

import io
import os
import sys
import json
import base64
import asyncio
import zipfile
import tempfile
from unittest.mock import patch

# Make ai_engine importable from repo root.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402


# --------------------------------------------------------------------------
# Formal spec mirror — isBugCondition (design "Formal Specification").
# Used as a PRECONDITION guard: every input these tests drive MUST be a
# non-bug input, so the preservation assertions are about the unchanged domain.
# --------------------------------------------------------------------------
_VISUAL_ROLES = {"cover", "content", "visual"}


def is_bug_condition(state: dict) -> bool:
    gate_suppressed = (
        state["htmlEnabled"]
        and state["vertexEnabled"]
        and state["role"] in _VISUAL_ROLES
        and not state["hasVertexImage"]
    )
    embed_discarded = (
        state["hasVertexImage"]
        and state["hasNativeDiagram"]
        and not state["hasImageFile"]
        and not state["hasSlideBg"]
    )
    return gate_suppressed or embed_discarded


# --------------------------------------------------------------------------
# Hermetic fakes
# --------------------------------------------------------------------------
_SECTION_TAG = 800_001          # bytes the section-HTML fake writes (slide #1)
_HTML_PNG_TAG = 900_000         # bytes the cover/_render_html_slide_to_png fake writes


def _make_png(tag: int) -> bytes:
    """Produce a unique, valid PNG so each generated image is byte-distinct."""
    img = Image.new(
        "RGB",
        (40 + (tag % 11), 30 + (tag % 7)),
        (tag % 256, (tag * 7) % 256, (tag * 31) % 256),
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class _FakeVertexClient:
    """Vertex client stub. ``enabled`` and ``fail`` are configurable so we can
    exercise the disabled path and the generate-failure path (Req 3.3)."""

    def __init__(self, enabled: bool = True, fail: bool = False) -> None:
        self.enabled = enabled
        self.fail = fail
        self.calls = 0
        self.generated_raw: list[bytes] = []

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.calls += 1
        if self.fail:
            return {"error": "quota-exceeded (test)"}
        raw = _make_png(self.calls)
        self.generated_raw.append(raw)
        return {"images": [base64.b64encode(raw).decode("ascii")]}


async def _img_gen_disabled(*_a, **_k):
    """Stand-in for _tool_generate_image — never returns a path (no network)."""
    return json.dumps({"error": "test-disabled"})


async def _render_html_png_fake(html, output_path, width=1920, height=1080, timeout=30, **_k):
    """Stand-in for _render_html_slide_to_png — writes a small valid PNG."""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "wb") as f:
        f.write(_make_png(_HTML_PNG_TAG))
    return {"ok": True}


def _make_section_html_fake(project_path: str, captured_profiles: list):
    """Build a stand-in for _generate_html_slide_for_section that records the
    style_profile it receives and writes a distinct PNG per call."""
    state = {"n": 0}

    async def _section_fake(gw, model_id, heading, body, doc_context,
                            proj_path, style_profile=None,
                            hero_image="", render_info=None):
        state["n"] += 1
        captured_profiles.append(style_profile)
        # task 3.5 — production now passes hero_image/render_info (additive,
        # byte-compatible). Mirror the contract: report layout/composited so
        # the caller's loss-zero branch behaves as in production. With Vertex
        # disabled/empty (non-bug preservation inputs) hero_image is "" so no
        # compositing occurs and the full-bleed slideBackground path is taken.
        if isinstance(render_info, dict):
            render_info["layout"] = "two_column"
            render_info["composited"] = bool(hero_image)
        gen = os.path.join(project_path, ".generated")
        os.makedirs(gen, exist_ok=True)
        name = f"slide-html-{state['n']}.png"
        with open(os.path.join(gen, name), "wb") as f:
            f.write(_make_png(_SECTION_TAG + state["n"] - 1))
        return f".generated/{name}"

    return _section_fake


# --------------------------------------------------------------------------
# PPTX inspection helpers
# --------------------------------------------------------------------------
def _pptx_media_bytes(pptx_abs: str) -> list[bytes]:
    with zipfile.ZipFile(pptx_abs) as z:
        return [z.read(n) for n in z.namelist() if n.startswith("ppt/media/")]


def _content_slide(pptx_abs: str):
    """Return the first content slide (index 1; index 0 is the cover)."""
    prs = Presentation(pptx_abs)
    assert len(prs.slides) >= 2, "expected cover + at least one content slide"
    return prs.slides[1]


def _shape_texts(slide) -> list[tuple[str, str]]:
    out = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append((str(sh.shape_type), sh.text_frame.text.strip()))
        except Exception:
            continue
    return out


def _fullbleed_pictures(slide) -> list:
    """PICTURE shapes that cover the whole 13.333×7.5in slide from (0,0)."""
    from pptx.util import Inches
    pics = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        if (sh.left == Inches(0) and sh.top == Inches(0)
                and abs(sh.width - Inches(13.333)) <= Inches(0.05)
                and abs(sh.height - Inches(7.5)) <= Inches(0.05)):
            pics.append(sh)
    return pics


# --------------------------------------------------------------------------
# Drivers — exercise the real _tool_generate_pptx decision seam
# --------------------------------------------------------------------------
def _run_native_html_off(slide: dict, *, vertex_enabled: bool, vertex_fail: bool = False,
                         style_profile=None):
    """HTML OFF driver (native / Vertex-embed path). Returns (fake, result)."""
    fake = _FakeVertexClient(enabled=vertex_enabled, fail=vertex_fail)
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",        # force HTML off
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex pre-gen eligible
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # skip LLM structuring + card fallback
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    tool_input = {"title": "T", "slides": [slide]}
    if style_profile is not None:
        tool_input["styleProfile"] = style_profile
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(tool_input, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return fake, result


def _run_html_on(slide: dict, *, vertex_enabled: bool, style_profile=None):
    """HTML ON driver (full-bleed path). Vertex disabled keeps the input non-bug
    (no gateSuppressed). Returns (fake, captured_profiles, result)."""
    fake = _FakeVertexClient(enabled=vertex_enabled)
    captured: list = []
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "1",        # force HTML on
        "AE_DISABLE_HTML_SLIDES": "0",
        # legacy(하이브리드 OFF) HTML 풀블리드 경로를 검증하는 드라이버.
        # 기본 ON(A안)으로 계약이 반전됐으므로 legacy 경로는 명시적 킬스위치로만 도달한다.
        "AE_HYBRID_RENDER": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    tool_input = {"title": "T", "slides": [slide]}
    if style_profile is not None:
        tool_input["styleProfile"] = style_profile
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", lambda *a, **k: object()), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_render_html_slide_to_png", _render_html_png_fake), \
            patch.object(srv, "_generate_html_slide_for_section",
                         _make_section_html_fake(proj, captured)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(tool_input, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return fake, captured, result


# --------------------------------------------------------------------------
# Input pools (curated so classification is deterministic)
# --------------------------------------------------------------------------
# Structural slides — _classify_section_diagram MUST return a real structural
# kind (flow / tree / architecture) for these (verified by observation).
_STRUCTURAL_SLIDES = [
    ("업무 처리 프로세스", ["접수", "검토", "승인", "완료"]),
    ("데이터 처리 흐름", ["수집", "정제", "분석", "시각화"]),
    ("처리 단계", ["요청 접수", "데이터 검증", "결과 반환"]),
    ("시스템 아키텍처", ["프론트엔드 계층", "백엔드 계층", "데이터 계층"]),
    ("배포 흐름 프로세스", ["빌드", "테스트", "스테이징", "배포"]),
]

# Non-structural content slides — _classify_section_diagram MUST return "" so the
# slide takes the plain content / Vertex-image path (NOT a native diagram).
_CONTENT_SLIDES = [
    ("환영합니다", ["반갑습니다", "함께 잘 부탁드립니다"]),
    ("팀 인사", ["좋은 하루입니다", "즐겁게 시작해요"]),
    ("오늘의 메시지", ["감사합니다", "끝까지 함께 가요"]),
    ("마무리 인사", ["수고하셨습니다", "다음에 또 만나요"]),
]

_STYLE_PROFILES = [
    {"primary": "#123456"},
    {"primary": "#0a7", "secondary": "#222", "font": "Pretendard"},
    {"palette": {"accent": "#ff6600"}, "primary": "#003366"},
]


# ==========================================================================
# PRES-1 (Req 3.1) — structural slides render as editable native shapes
# ==========================================================================
@settings(max_examples=10, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES))
def test_pres1_structural_renders_native_shapes(slide):
    title, bullets = slide

    # Precondition: this really classifies as a real structural diagram.
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind in ("flow", "tree", "architecture"), (
        f"테스트 전제 실패: {title!r}는 구조형으로 분류돼야 함 (kind={kind!r})"
    )
    # Precondition: non-bug input (Vertex skipped for structural → no image;
    # HTML off → no gateSuppressed).
    assert not is_bug_condition({
        "htmlEnabled": False, "vertexEnabled": True, "role": "structural",
        "hasVertexImage": False, "hasNativeDiagram": True,
        "hasImageFile": False, "hasSlideBg": False,
    })

    fake, result = _run_native_html_off({"title": title, "bullets": list(bullets)},
                                        vertex_enabled=True)

    # (Req 3.1) Vertex pre-gen is skipped for structural slides — no raster image.
    assert fake.calls == 0, (
        f"구조형 슬라이드는 Vertex 래스터를 생성하지 않아야 함 (calls={fake.calls})"
    )
    assert _pptx_media_bytes(result["absPath"]) == [], (
        "구조형 슬라이드에는 래스터 이미지가 임베드되면 안 됨 (편집 가능 네이티브 도형이어야 함)"
    )
    # The bullet labels survive as EDITABLE native shape text (not rasterised).
    texts = " ".join(t for _, t in _shape_texts(_content_slide(result["absPath"])))
    hits = [b for b in bullets if b in texts]
    assert len(hits) >= max(2, len(bullets) - 1), (
        f"구조형 라벨이 편집 가능 도형 텍스트로 남아야 함 — hits={hits}, bullets={bullets}"
    )


# ==========================================================================
# PRES-2 (Req 3.2) — HTML full-bleed slideBackground path unchanged
# ==========================================================================
@settings(max_examples=8, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES))
def test_pres2_html_fullbleed_path_preserved(slide):
    title, bullets = slide

    # Precondition: non-structural content (so HTML renders it, not a diagram).
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind == "", f"테스트 전제 실패: {title!r}는 비구조형이어야 함 (kind={kind!r})"
    # Precondition: HTML on + Vertex OFF → non-bug (no gateSuppressed).
    assert not is_bug_condition({
        "htmlEnabled": True, "vertexEnabled": False, "role": "content",
        "hasVertexImage": False, "hasNativeDiagram": False,
        "hasImageFile": False, "hasSlideBg": False,
    })

    fake, captured, result = _run_html_on(
        {"title": title, "bullets": list(bullets)}, vertex_enabled=False)

    # Vertex disabled → never called.
    assert fake.calls == 0, f"Vertex 비활성 시 호출되면 안 됨 (calls={fake.calls})"
    # The section HTML capture is embedded as a (0,0) full-bleed background.
    content = _content_slide(result["absPath"])
    fb = _fullbleed_pictures(content)
    assert fb, "HTML 풀블리드 경로: 콘텐츠 슬라이드에 (0,0) 풀블리드 PICTURE가 있어야 함"
    media = _pptx_media_bytes(result["absPath"])
    assert _make_png(_SECTION_TAG) in media, (
        "slideBackground 경로: 섹션 HTML→PNG 캡처 바이트가 임베드돼야 함"
    )


# ==========================================================================
# PRES-3 (Req 3.3) — Vertex disabled / failed → native fallback preserved
# ==========================================================================
@settings(max_examples=12, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       fail_mode=st.sampled_from(["disabled", "generate_fails"]))
def test_pres3_vertex_unavailable_native_fallback(slide, fail_mode):
    title, bullets = slide

    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind == "", f"테스트 전제 실패: {title!r}는 비구조형이어야 함 (kind={kind!r})"
    # Precondition: Vertex effectively unavailable → no Vertex image → non-bug.
    assert not is_bug_condition({
        "htmlEnabled": False, "vertexEnabled": (fail_mode != "disabled"),
        "role": "content", "hasVertexImage": False, "hasNativeDiagram": False,
        "hasImageFile": False, "hasSlideBg": False,
    })

    enabled = fail_mode != "disabled"
    fail = fail_mode == "generate_fails"
    fake, result = _run_native_html_off(
        {"title": title, "bullets": list(bullets)},
        vertex_enabled=enabled, vertex_fail=fail)

    # (Req 3.3) No Vertex image ends up embedded — native text fallback.
    assert _pptx_media_bytes(result["absPath"]) == [], (
        f"Vertex 미가용({fail_mode}) 시 래스터가 임베드되면 안 됨 — 네이티브 폴백이어야 함"
    )
    # The body bullet text is preserved on the content slide.
    texts = " ".join(t for _, t in _shape_texts(_content_slide(result["absPath"])))
    hits = [b for b in bullets if b in texts]
    assert hits, (
        f"네이티브 폴백 시 본문 불릿 텍스트가 보존돼야 함 — bullets={bullets}, texts={texts!r}"
    )


# ==========================================================================
# PRES-4 (Req 3.5) — styleProfile inheritance (HTML token + native palette)
# ==========================================================================
@settings(max_examples=9, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       profile=st.sampled_from(_STYLE_PROFILES))
def test_pres4_style_profile_inheritance_preserved(slide, profile):
    title, bullets = slide

    # HTML token inheritance: the per-section renderer receives the exact profile.
    fake, captured, result = _run_html_on(
        {"title": title, "bullets": list(bullets)},
        vertex_enabled=False, style_profile=profile)
    assert profile in captured, (
        f"styleProfile가 HTML 섹션 렌더러로 전달돼야 함(토큰 상속) — captured={captured}"
    )

    # Native palette inheritance: _build_palette derives a stable palette from the
    # same profile (the path native diagrams use). Must not raise; deterministic.
    palette_a = srv._build_palette(profile)
    palette_b = srv._build_palette(profile)
    assert palette_a == palette_b, "동일 styleProfile은 동일 네이티브 팔레트를 내야 함(결정성)"


# ==========================================================================
# PRES-5 — caller-specified imageFile / slideBackground precedence preserved
# ==========================================================================
def _write_wide_png(path: str, w: int = 1200, h: int = 700) -> bytes:
    img = Image.new("RGB", (w, h), (11, 22, 33))
    img.save(path, format="PNG")
    with open(path, "rb") as f:
        return f.read()


def _write_bg_png(path: str) -> bytes:
    img = Image.new("RGB", (1920, 1080), (7, 8, 9))
    img.save(path, format="PNG")
    with open(path, "rb") as f:
        return f.read()


@settings(max_examples=8, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       which=st.sampled_from(["imageFile", "slideBackground"]))
def test_pres5_caller_specified_image_precedence(slide, which):
    title, bullets = slide

    fake = _FakeVertexClient(enabled=True)
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex eligible — but must skip this slide
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    png_path = os.path.join(proj, "caller.png")
    if which == "imageFile":
        caller_bytes = _write_wide_png(png_path)
    else:
        caller_bytes = _write_bg_png(png_path)
    sd = {"title": title, "bullets": list(bullets), which: png_path}

    # Precondition: caller image present → non-bug (embedDiscarded false).
    assert not is_bug_condition({
        "htmlEnabled": False, "vertexEnabled": True, "role": "content",
        "hasVertexImage": False, "hasNativeDiagram": False,
        "hasImageFile": which == "imageFile", "hasSlideBg": which == "slideBackground",
    })

    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": [sd]},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"

    # Vertex skips a slide that already has imageFile/slideBackground.
    assert fake.calls == 0, (
        f"caller가 {which}를 지정한 슬라이드는 Vertex 생성을 건너뛰어야 함 (calls={fake.calls})"
    )
    # The caller's exact image bytes are the ones embedded (precedence preserved).
    media = _pptx_media_bytes(result["absPath"])
    assert caller_bytes in media, (
        f"caller가 지정한 {which} 이미지가 그대로 임베드돼야 함(우선순위 보존)"
    )


# ==========================================================================
# PROPERTY 5 — Gateway Constraint (이미지 외 호출은 Vertex 미사용)
#
# Design "Property 5: 게이트웨이 제약 보존 — 이미지 외 호출은 Vertex 미사용":
#   _For any_ input, LLM / operation-JSON generation flows ONLY through the
#   Bedrock Gateway, and Vertex is invoked ONLY on the image-generation path
#   (zero Vertex calls for non-image / non-visual work).  gateway.md image
#   exception clause: "LLM/추론/operation JSON 생성: Bedrock Gateway 경유 그대로
#   유지 (예외 없음)", "Vertex 호출은 ... 이미지 생성 경로에서만".
#
# Hermetic mock-spy strategy:
#   * ``_SpyVertexClient`` is a STRICT spy — it records the ``model_class`` of
#     every ``generate()`` call and, via ``__getattr__``, records ANY access to
#     a member other than the legitimate image-client surface (``enabled`` /
#     ``generate``).  If the production pptx path ever tried to repurpose the
#     Vertex client as an LLM / operation generator (e.g. ``.converse`` /
#     ``.chat`` / ``.complete``), that access is recorded AND raises — proving
#     structurally that Vertex is confined to the image API.
#   * ``_GwSpy`` records every Bedrock-Gateway acquisition (``_get_gw``) so we
#     can assert the gateway IS the seam used for LLM / operation work.
#
# _Requirements: 3.4_
# ==========================================================================
_LLM_LIKE_MEMBERS = {
    "converse", "chat", "complete", "completion", "stream", "invoke",
    "invoke_model", "responses", "openai_responses_sync", "generate_text",
    "messages", "embed",
}


class _SpyVertexClient:
    """Strict Vertex spy for Property 5.

    Records every ``generate()`` call's ``model_class`` and flags any access to
    an unexpected member (anything that is not ``enabled`` / ``generate``). The
    production pptx path only ever touches ``.enabled`` and ``.generate(...)``;
    any other access — especially an LLM-like verb — is captured in
    ``foreign_access`` and also raises ``AttributeError`` so a mis-route surfaces
    immediately.
    """

    def __init__(self, enabled: bool = True) -> None:
        # set the recording lists FIRST so __getattr__ can reference them safely
        self.foreign_access: list[str] = []
        self.generate_calls: list = []          # model_class per generate() call
        self.enabled = enabled

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.generate_calls.append(model_class)
        raw = _make_png(900 + len(self.generate_calls))
        return {"images": [base64.b64encode(raw).decode("ascii")]}

    @property
    def llm_like_access(self) -> list:
        return [n for n in self.foreign_access if n in _LLM_LIKE_MEMBERS]

    def __getattr__(self, name: str):
        # Reached only for members not found normally (i.e. not enabled/generate
        # /the recording lists). Record then mimic a normal missing attribute.
        if name.startswith("__"):
            raise AttributeError(name)
        object.__getattribute__(self, "foreign_access").append(name)
        raise AttributeError(
            f"Vertex client accessed unexpected member {name!r} — Vertex must "
            "only be used via generate() on the image-generation path"
        )


class _GwSpy:
    """Records Bedrock-Gateway acquisitions so we can assert the gateway is the
    LLM / operation seam. Returns an opaque object (the real gw is never hit)."""

    def __init__(self) -> None:
        self.calls = 0

    def __call__(self, *_a, **_k):
        self.calls += 1
        return object()


def _run_pptx_html_off_with(client, slides, *, env_extra=None):
    """HTML-OFF driver that injects a caller-supplied Vertex ``client`` (spy)."""
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    if env_extra:
        env.update(env_extra)
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: client), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": list(slides)},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return result


def _run_pptx_html_on_with(client, slides, *, gw_factory):
    """HTML-ON driver that injects a Vertex ``client`` (spy) and a gateway
    factory spy (``gw_factory``) so the gateway seam is observable."""
    proj = tempfile.mkdtemp()
    captured: list = []
    env = {
        "AE_ENABLE_HTML_SLIDES": "1",
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: client), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", gw_factory), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_render_html_slide_to_png", _render_html_png_fake), \
            patch.object(srv, "_generate_html_slide_for_section",
                         _make_section_html_fake(proj, captured)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": list(slides)},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return result


# Visual-intent slides — non-structural title/bullets + a photographic/illustrative
# imagePrompt (no structural signals) so _classify_slide_role => "visual" and the
# Vertex image-generation path IS exercised.
_VISUAL_SLIDES = [
    {"title": "회사 소개", "bullets": ["신뢰", "혁신"],
     "imagePrompt": "a high quality professional photograph of a modern corporate office, natural light"},
    {"title": "브랜드 비전", "bullets": ["미래 지향"],
     "imagePrompt": "a cinematic photograph of a sunrise over a city skyline, warm tones"},
    {"title": "고객 경험", "bullets": ["만족"],
     "imagePrompt": "an editorial illustration of happy customers using a product, flat design"},
    {"title": "팀 문화", "bullets": ["협업"],
     "imagePrompt": "a candid photograph of a diverse team collaborating in a bright studio"},
]


# --------------------------------------------------------------------------
# PROP5-1 (Req 3.4) — non-visual (structural) work invokes ZERO Vertex calls
# --------------------------------------------------------------------------
@settings(max_examples=10, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES))
def test_prop5_nonvisual_structural_uses_no_vertex(slide):
    title, bullets = slide

    # Precondition: genuinely structural (non-visual) work.
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind in ("flow", "tree", "architecture"), (
        f"테스트 전제 실패: {title!r}는 구조형(비주얼 아님)이어야 함 (kind={kind!r})"
    )
    role = srv._classify_slide_role({"title": title, "bullets": list(bullets)}, False, "doc")
    assert role == "structural", f"전제 실패: role={role!r} (structural 기대)"

    spy = _SpyVertexClient(enabled=True)
    _run_pptx_html_off_with(spy, [{"title": title, "bullets": list(bullets)}])

    # (Req 3.4) Vertex is NOT called for non-image / non-visual work.
    assert spy.generate_calls == [], (
        f"이미지 외(구조형) 작업에서 Vertex가 호출되면 안 됨 — generate_calls={spy.generate_calls}"
    )
    # Vertex was never repurposed as an LLM / operation generator.
    assert spy.llm_like_access == [], (
        f"Vertex 클라이언트가 LLM/operation 용도로 접근되면 안 됨 — {spy.llm_like_access}"
    )
    assert spy.foreign_access == [], (
        f"Vertex 클라이언트는 enabled/generate 외 멤버가 접근되면 안 됨 — {spy.foreign_access}"
    )


# --------------------------------------------------------------------------
# PROP5-2 (Req 3.4) — Vertex is reached ONLY via the image-generation API
# --------------------------------------------------------------------------
@settings(max_examples=10, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_VISUAL_SLIDES))
def test_prop5_vertex_only_via_image_generation_path(slide):
    # Precondition: visual-intent, non-structural → role "visual" (image work).
    kind, _ = srv._classify_section_diagram(
        slide["title"], "\n".join(slide["bullets"]), "doc")
    assert kind == "", f"전제 실패: {slide['title']!r}는 비구조형이어야 함 (kind={kind!r})"
    assert srv._has_visual_intent(slide["imagePrompt"]), "전제 실패: visual intent 여야 함"
    role = srv._classify_slide_role(slide, False, "doc")
    assert role == "visual", f"전제 실패: role={role!r} (visual 기대)"

    spy = _SpyVertexClient(enabled=True)
    _run_pptx_html_off_with(spy, [dict(slide)])

    # Vertex WAS exercised (image work happened on the image path).
    assert spy.generate_calls, (
        "비주얼 슬라이드는 Vertex 이미지 생성 경로를 거쳐야 함 (generate 호출 0)"
    )
    # Every Vertex invocation used an IMAGE-generation model_class — never a
    # text / LLM class. This is what confines Vertex to the image path.
    for mc in spy.generate_calls:
        assert isinstance(mc, str) and "image" in mc.lower(), (
            f"Vertex 호출은 이미지 생성 model_class 여야 함 — got {mc!r}"
        )
    # Vertex was never used as an LLM / operation generator.
    assert spy.llm_like_access == [], (
        f"Vertex 클라이언트가 LLM/operation 용도로 접근되면 안 됨 — {spy.llm_like_access}"
    )
    assert spy.foreign_access == [], (
        f"Vertex 클라이언트는 enabled/generate 외 멤버가 접근되면 안 됨 — {spy.foreign_access}"
    )


# --------------------------------------------------------------------------
# PROP5-3 (Req 3.4) — LLM / operation work flows through the Bedrock Gateway,
#                     never through Vertex (even under the HTML quality path)
# --------------------------------------------------------------------------
@settings(max_examples=10, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES))
def test_prop5_llm_operation_via_gateway_not_vertex(slide):
    title, bullets = slide

    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind in ("flow", "tree", "architecture"), (
        f"전제 실패: {title!r}는 구조형(비주얼 아님)이어야 함 (kind={kind!r})"
    )

    spy = _SpyVertexClient(enabled=True)
    gw_spy = _GwSpy()
    _run_pptx_html_on_with(spy, [{"title": title, "bullets": list(bullets)}],
                           gw_factory=gw_spy)

    # The Bedrock Gateway IS the seam acquired for LLM / operation work.
    assert gw_spy.calls >= 1, (
        "LLM/operation 생성은 Bedrock Gateway(_get_gw) 경유여야 함 — 게이트웨이 미취득"
    )
    # Vertex is untouched for non-visual (structural) work, even on the HTML path.
    assert spy.generate_calls == [], (
        f"구조형(비주얼 아님) 작업에서 Vertex가 호출되면 안 됨 — generate_calls={spy.generate_calls}"
    )
    assert spy.llm_like_access == [], (
        f"Vertex 클라이언트가 LLM/operation 용도로 접근되면 안 됨 — {spy.llm_like_access}"
    )
    assert spy.foreign_access == [], (
        f"Vertex 클라이언트는 enabled/generate 외 멤버가 접근되면 안 됨 — {spy.foreign_access}"
    )


# ==========================================================================
# HYBRID-RENDER REGRESSION EXTENSIONS
# spec: pptx-ultra-quality-hybrid-render — Tasks 8.2 / 8.3 / 8.4.
#
# These extend the preservation patterns above (pres5 / pres3 / prop5) to the
# hybrid render routing seam, driven with the opt-in flag ``AE_HYBRID_RENDER=1``
# turned ON. They prove that enabling hybrid routing does NOT regress the R4
# preservation guarantees: caller-media precedence (Property 15), Vertex
# unavailable/failed loss-zero fallback (Property 16), and the gateway
# constraint that Vertex is never used for non-image work (Property 17).
#
# Everything remains hermetic — Vertex (``get_vertex_image_client`` /
# ``generate``), the Bedrock Gateway (``_get_gw``) and the HTML→PNG renderer
# are mocked. No network. Each property test runs 100+ hypothesis iterations.
# ==========================================================================
_HYBRID_ON = {"AE_HYBRID_RENDER": "1"}


# ==========================================================================
# Feature: pptx-ultra-quality-hybrid-render, Property 15: caller 지정 미디어
# 우선순위 보존 — caller가 imageFile/slideBackground를 지정한 슬라이드는 하이브리드
# 라우팅이 덮어쓰지 않고 기존 _select_render_plan에 위임되며 Vertex 사전생성이 스킵된다.
# (기존 pres5 확장 — AE_HYBRID_RENDER=1 ON)
# Validates: Requirements 4.2
# ==========================================================================
def _run_caller_precedence_hybrid(slide_dict: dict):
    """HTML-OFF driver with AE_HYBRID_RENDER=1 ON. Injects a Vertex spy so we can
    assert pre-gen is SKIPPED for a caller-media slide even under hybrid routing.
    Returns (fake, result)."""
    fake = _FakeVertexClient(enabled=True)
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex eligible — must still skip this slide
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
        **_HYBRID_ON,                        # hybrid routing ON
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": [slide_dict]},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return fake, result


@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       which=st.sampled_from(["imageFile", "slideBackground"]),
       salt=st.integers(min_value=0, max_value=10 ** 9))
def test_prop15_caller_media_precedence_under_hybrid(slide, which, salt):
    # salt widens the input space (classification-stable title suffix) so the
    # property runs 100+ distinct hypothesis iterations.
    base_title, bullets = slide
    title = f"{base_title} #{salt}"

    proj = tempfile.mkdtemp()
    png_path = os.path.join(proj, "caller.png")
    caller_bytes = (_write_wide_png(png_path) if which == "imageFile"
                    else _write_bg_png(png_path))
    sd = {"title": title, "bullets": list(bullets), which: png_path}

    # Precondition: caller media present → non-bug input (embedDiscarded false),
    # and this is the exact domain Property 15 governs (caller precedence).
    assert not is_bug_condition({
        "htmlEnabled": False, "vertexEnabled": True, "role": "content",
        "hasVertexImage": False, "hasNativeDiagram": False,
        "hasImageFile": which == "imageFile", "hasSlideBg": which == "slideBackground",
    })

    fake, result = _run_caller_precedence_hybrid(sd)

    # (R4.2) Hybrid routing did NOT trigger Vertex pre-gen for a caller-media slide.
    assert fake.calls == 0, (
        f"하이브리드 ON에서도 caller가 {which}를 지정한 슬라이드는 Vertex 생성을 "
        f"건너뛰어야 함 (calls={fake.calls})"
    )
    # (R4.2) The caller's exact bytes are the ones embedded — routing did not
    # overwrite the caller media (delegated to existing _select_render_plan).
    media = _pptx_media_bytes(result["absPath"])
    assert caller_bytes in media, (
        f"하이브리드 라우팅이 caller 지정 {which} 이미지를 덮어쓰면 안 됨(우선순위 보존)"
    )


# ==========================================================================
# Feature: pptx-ultra-quality-hybrid-render, Property 16: Vertex 비활성/실패
# 폴백의 손실-0 — vertex_enabled==False 또는 generate 실패 시, 렌더는 콘텐츠 손실
# 항목 개수 0으로 편집 가능 네이티브/HTML 폴백으로 전환하고 폴백 발생을 표시한다.
# (기존 pres3 확장 — AE_HYBRID_RENDER=1 ON)
# Validates: Requirements 4.3, 6.3
# ==========================================================================
def _run_native_html_off_hybrid(slide: dict, *, vertex_enabled: bool, vertex_fail: bool):
    """pres3-style HTML-OFF native driver, but with AE_HYBRID_RENDER=1 ON."""
    fake = _FakeVertexClient(enabled=vertex_enabled, fail=vertex_fail)
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
        **_HYBRID_ON,
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx({"title": "T", "slides": [slide]},
                                                  project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    return fake, result


@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_CONTENT_SLIDES),
       fail_mode=st.sampled_from(["disabled", "generate_fails"]),
       salt=st.integers(min_value=0, max_value=10 ** 9))
def test_prop16_vertex_unavailable_lossless_fallback_under_hybrid(slide, fail_mode, salt):
    # salt widens the input space (classification-stable title suffix) so the
    # property runs 100+ distinct hypothesis iterations.
    base_title, bullets = slide
    title = f"{base_title} #{salt}"

    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind == "", f"테스트 전제 실패: {title!r}는 비구조형이어야 함 (kind={kind!r})"
    # Precondition: Vertex effectively unavailable → no Vertex image → non-bug.
    assert not is_bug_condition({
        "htmlEnabled": False, "vertexEnabled": (fail_mode != "disabled"),
        "role": "content", "hasVertexImage": False, "hasNativeDiagram": False,
        "hasImageFile": False, "hasSlideBg": False,
    })

    enabled = fail_mode != "disabled"
    fail = fail_mode == "generate_fails"
    fake, result = _run_native_html_off_hybrid(
        {"title": title, "bullets": list(bullets)},
        vertex_enabled=enabled, vertex_fail=fail)

    content = _content_slide(result["absPath"])
    texts = " ".join(t for _, t in _shape_texts(content))

    # (R4.3/6.3) Content-loss item count == 0: every caller bullet survives as
    # EDITABLE shape text AND no raster image was embedded (native fallback).
    lost_bullets = [b for b in bullets if b not in texts]
    embedded_media = _pptx_media_bytes(result["absPath"])
    content_loss_items = len(lost_bullets) + len(embedded_media)
    assert content_loss_items == 0, (
        f"Vertex 미가용({fail_mode}) 폴백에서 콘텐츠 손실 항목 개수는 0이어야 함 — "
        f"lost_bullets={lost_bullets}, embedded_media={len(embedded_media)}"
    )
    # (R4.3/6.3) The fallback is editable: text runs exist (not rasterised) and
    # NO full-bleed PICTURE was baked in.
    assert texts.strip(), "폴백 슬라이드는 편집 가능 텍스트 run을 보유해야 함(래스터화 아님)"
    assert _fullbleed_pictures(content) == [], (
        "폴백 경로는 풀블리드 PICTURE를 굽지 않아야 함(편집 가능 네이티브/HTML)"
    )
    # (R6.3) Fallback occurrence is indicated by the observable Vertex-call
    # signature: disabled → never called; generate_fails → attempted then fell
    # back. In both cases NO Vertex raster reached the deck (loss-0 confirmed).
    if fail_mode == "disabled":
        assert fake.calls == 0, "Vertex 비활성 시 호출 0이어야 함(폴백 표시)"
    else:
        assert fake.calls >= 1, (
            "generate 실패 모드는 Vertex 시도 후 폴백해야 함(폴백 발생 표시)"
        )


# ==========================================================================
# Feature: pptx-ultra-quality-hybrid-render, Property 17: 게이트웨이 제약 —
# 이미지 외 Vertex 미호출 — 이미지 생성 경로가 아닌 실행(라우팅 결정, 프롬프트 빌드,
# operation JSON 생성)에서 VertexImageClient 호출 개수는 0이고 모든 LLM 호출은
# Bedrock_Gateway(_get_gw) 경유로만 발생한다. (기존 prop5 확장 — AE_HYBRID_RENDER=1 ON)
# Validates: Requirements 3.7, 4.4, 6.3
# ==========================================================================
def _run_pptx_html_off_with_hybrid(client, slides):
    """prop5-style HTML-OFF spy driver, hybrid ON (reuses env_extra seam)."""
    return _run_pptx_html_off_with(client, slides, env_extra=dict(_HYBRID_ON))


def _run_pptx_html_on_with_hybrid(client, slides, *, gw_factory):
    """prop5-style HTML-ON spy driver, hybrid ON. Wraps the existing helper with
    an outer patch.dict so AE_HYBRID_RENDER=1 is set for the whole run."""
    with patch.dict(os.environ, dict(_HYBRID_ON), clear=False):
        return _run_pptx_html_on_with(client, slides, gw_factory=gw_factory)


@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.function_scoped_fixture, HealthCheck.too_slow])
@given(slide=st.sampled_from(_STRUCTURAL_SLIDES), html_on=st.booleans(),
       salt=st.integers(min_value=0, max_value=10 ** 9))
def test_prop17_gateway_constraint_non_image_no_vertex_under_hybrid(slide, html_on, salt):
    # salt widens the input space (classification-stable title suffix) so the
    # property runs 100+ distinct hypothesis iterations.
    base_title, bullets = slide
    title = f"{base_title} #{salt}"

    # Precondition: genuinely structural (non-image) work — routing decision,
    # prompt build and operation JSON generation, but NOT the image path.
    kind, _ = srv._classify_section_diagram(title, "\n".join(bullets), "doc")
    assert kind in ("flow", "tree", "architecture"), (
        f"전제 실패: {title!r}는 구조형(비주얼 아님)이어야 함 (kind={kind!r})"
    )
    role = srv._classify_slide_role({"title": title, "bullets": list(bullets)}, False, "doc")
    assert role == "structural", f"전제 실패: role={role!r} (structural 기대)"

    # Sanity: the pure routing/prompt seams themselves never touch a Vertex
    # client (they take no client and construct none) — hybrid plan for a
    # structural role never routes to a full-bleed Vertex primary.
    plan = srv._select_hybrid_render_plan(
        role="structural", vertex_enabled=True, html_enabled=html_on,
        has_vertex_image=False, has_native_diagram=True,
        has_image_file=False, has_slide_bg=False)
    assert plan["primary"] == "NATIVE_SHAPES", (
        f"structural 역할의 하이브리드 primary는 NATIVE_SHAPES 여야 함 — {plan}"
    )

    spy = _SpyVertexClient(enabled=True)
    if html_on:
        gw_spy = _GwSpy()
        _run_pptx_html_on_with_hybrid(
            spy, [{"title": title, "bullets": list(bullets)}], gw_factory=gw_spy)
        # (R4.4/6.3) LLM / operation work flows through the Bedrock Gateway.
        assert gw_spy.calls >= 1, (
            "LLM/operation 생성은 Bedrock Gateway(_get_gw) 경유여야 함 — 게이트웨이 미취득"
        )
    else:
        _run_pptx_html_off_with_hybrid(
            spy, [{"title": title, "bullets": list(bullets)}])

    # (R3.7/4.4) Non-image work → ZERO Vertex calls under hybrid routing.
    assert spy.generate_calls == [], (
        f"이미지 외(구조형) 작업에서 Vertex가 호출되면 안 됨 — generate_calls={spy.generate_calls}"
    )
    # Vertex was never repurposed as an LLM / operation generator.
    assert spy.llm_like_access == [], (
        f"Vertex 클라이언트가 LLM/operation 용도로 접근되면 안 됨 — {spy.llm_like_access}"
    )
    assert spy.foreign_access == [], (
        f"Vertex 클라이언트는 enabled/generate 외 멤버가 접근되면 안 됨 — {spy.foreign_access}"
    )


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
