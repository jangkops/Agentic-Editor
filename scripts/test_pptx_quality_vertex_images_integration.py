"""Hermetic integration test — spec: pptx-quality-vertex-images (bugfix), Task 7.

Drives the REAL ``ai_engine/server.py:_tool_generate_pptx`` decision seam end to
end over a MIXED deck (cover + high-density content + structural + photo/visual)
and asserts the spec's three integration guarantees (design "Integration Tests"):

  (a) every Vertex image that is generated ends up embedded — zero
      "generated but unused" (loss-zero, design Property 1 / 3);
  (b) genuine structural slides (flow / tree / architecture) stay EDITABLE
      native shapes — no rasterised full-bleed image over them (Req 3.1);
  (c) the produced PPTX slide count and background/visual embedding match
      expectation (cover + one slide per input section).

Two complementary scenarios cover the mixed deck:

  SCENARIO 1 — HTML OFF (``test_integration_mixed_deck_html_off``)
      Vertex images are embedded as REAL bytes (visual or full-bleed backdrop),
      so loss-zero is verified at the byte level: every byte sequence the mocked
      Vertex client produced is found inside ``ppt/media/*``. The structural
      slide is skipped by Vertex pre-gen and renders as native text shapes.

  SCENARIO 2 — HTML ON (``test_integration_mixed_deck_html_on_coexistence``)
      HTML high-density layouts are the primary renderer AND Vertex coexists
      (the fix removed the ``not _html_enabled`` mutual-exclusion gate). The
      section renderer fake bakes the Vertex hero bytes into the composited
      slide-background PNG, so the generated Vertex bytes still surface in
      ``ppt/media`` — proving the image is not discarded when HTML is on.

Everything is hermetic — no network. The Bedrock gateway, the Vertex client
(``get_vertex_image_client`` / ``generate``), the HTML→PNG renderer and
``_tool_generate_image`` are all mocked. Helpers mirror the existing
``test_pptx_quality_vertex_images_*`` test files.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_quality_vertex_images_integration.py -p no:cacheprovider -q

_Requirements: 2.1, 2.2, 2.3, 2.4, 3.1, 3.2_
"""
from __future__ import annotations

import io
import os
import sys
import json
import base64
import asyncio
import zipfile
import tempfile
from unittest.mock import patch

# Make ai_engine importable from repo root.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402


# --------------------------------------------------------------------------
# Hermetic fakes (mirror the existing spec test files)
# --------------------------------------------------------------------------
def _make_png(tag: int) -> bytes:
    """Produce a unique, valid PNG so each generated image is byte-distinct."""
    img = Image.new(
        "RGB",
        (40 + (tag % 11), 30 + (tag % 7)),
        (tag % 256, (tag * 7) % 256, (tag * 31) % 256),
    )
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class _FakeVertexClient:
    """Stand-in for VertexImageClient — always 'enabled', returns unique PNGs.

    Each ``generate`` call writes a byte-distinct PNG and records the raw bytes
    so the test can verify every generated image is embedded (loss-zero)."""

    def __init__(self) -> None:
        self.enabled = True
        self.calls = 0
        self.generated_raw: list[bytes] = []

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.calls += 1
        raw = _make_png(1000 + self.calls)
        self.generated_raw.append(raw)
        return {"images": [base64.b64encode(raw).decode("ascii")]}


async def _img_gen_disabled(*_a, **_k):
    """Stand-in for _tool_generate_image — never returns a path (no network)."""
    return json.dumps({"error": "test-disabled"})


async def _render_html_png_fake(html, output_path, width=1920, height=1080, timeout=30, **_k):
    """Stand-in for _render_html_slide_to_png (cover) — writes a valid PNG."""
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    with open(output_path, "wb") as f:
        f.write(_make_png(900_000))
    return {"ok": True}


def _make_compositing_section_fake(project_path: str, captured_heroes: list):
    """Section-HTML fake that COMPOSITES the Vertex hero (task 3.5).

    When ``hero_image`` is supplied it reads those exact Vertex bytes and writes
    them as the rendered section PNG, then reports ``composited=True`` — faithful
    to production, where the composited single PNG becomes the slideBackground.
    With no hero it writes its own distinct PNG. Records every hero path seen."""
    state = {"n": 0}

    async def _section_fake(gw, model_id, heading, body, doc_context,
                            proj_path, style_profile=None,
                            hero_image="", render_info=None):
        state["n"] += 1
        captured_heroes.append(hero_image)
        gen = os.path.join(project_path, ".generated")
        os.makedirs(gen, exist_ok=True)
        name = f"slide-html-{state['n']}.png"
        out = os.path.join(gen, name)
        if hero_image:
            # Bake the Vertex hero bytes into the composited section PNG.
            hero_abs = hero_image if os.path.isabs(hero_image) \
                else os.path.join(project_path, hero_image)
            with open(hero_abs, "rb") as hf:
                data = hf.read()
            with open(out, "wb") as f:
                f.write(data)
            if isinstance(render_info, dict):
                render_info["layout"] = "two_column"
                render_info["composited"] = True
        else:
            with open(out, "wb") as f:
                f.write(_make_png(800_000 + state["n"]))
            if isinstance(render_info, dict):
                render_info["layout"] = "feature_grid"
                render_info["composited"] = False
        return f".generated/{name}"

    return _section_fake


# --------------------------------------------------------------------------
# PPTX inspection helpers
# --------------------------------------------------------------------------
def _pptx_media_bytes(pptx_abs: str) -> list[bytes]:
    with zipfile.ZipFile(pptx_abs) as z:
        return [z.read(n) for n in z.namelist() if n.startswith("ppt/media/")]


def _vertex_usage(fake: _FakeVertexClient, pptx_abs: str):
    """Return (generated, embedded, unused) counts for the Vertex images."""
    media = _pptx_media_bytes(pptx_abs)
    embedded = sum(1 for raw in fake.generated_raw if raw in media)
    generated = len(fake.generated_raw)
    return generated, embedded, generated - embedded


def _slides(pptx_abs: str):
    return list(Presentation(pptx_abs).slides)


def _shape_texts(slide) -> list[str]:
    out = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                out.append(sh.text_frame.text.strip())
        except Exception:
            continue
    return out


def _fullbleed_pictures(slide) -> list:
    """PICTURE shapes covering the whole 13.333×7.5in slide from (0,0)."""
    from pptx.util import Inches
    pics = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        if (sh.left == Inches(0) and sh.top == Inches(0)
                and abs(sh.width - Inches(13.333)) <= Inches(0.05)
                and abs(sh.height - Inches(7.5)) <= Inches(0.05)):
            pics.append(sh)
    return pics


# --------------------------------------------------------------------------
# Mixed deck — curated so classification is deterministic.
# Order: [structural(flow), content(high-density), visual(photo)]
# --------------------------------------------------------------------------
_STRUCTURAL = {"title": "업무 처리 프로세스", "bullets": ["접수", "검토", "승인", "완료"]}
_CONTENT = {"title": "핵심 기능 요약",
            "bullets": ["빠른 처리", "안정적 운영", "유연한 확장", "강력한 보안",
                        "비용 절감", "쉬운 사용성"]}
_VISUAL = {"title": "회사 소개",
           "bullets": ["신뢰를 최우선으로"],
           "imagePrompt": "a high quality professional photograph of a modern "
                          "corporate office, natural light, wide angle"}


def _assert_deck_classification():
    """Precondition: the curated deck classifies as intended."""
    k_struct, _ = srv._classify_section_diagram(
        _STRUCTURAL["title"], "\n".join(_STRUCTURAL["bullets"]), "doc")
    assert k_struct in ("flow", "tree", "architecture"), (
        f"전제 실패: 구조형 슬라이드가 structural로 분류돼야 함 (kind={k_struct!r})")
    assert srv._classify_slide_role(_STRUCTURAL, False, "doc") == "structural"

    # content 슬라이드: 진짜 구조형(flow/tree/architecture)이 아니므로 role=content 로
    # 흡수돼야 한다(kpi/cards/twocol 등 고밀도 콘텐츠는 structural 아님 — design §1).
    k_content, _ = srv._classify_section_diagram(
        _CONTENT["title"], "\n".join(_CONTENT["bullets"]), "doc")
    assert k_content not in ("flow", "tree", "architecture"), (
        f"전제 실패: 콘텐츠 슬라이드는 진짜 구조형이 아니어야 함 (kind={k_content!r})")
    assert srv._classify_slide_role(_CONTENT, False, "doc") == "content", (
        "전제 실패: 고밀도 콘텐츠 슬라이드는 content로 분류돼야 함")

    assert srv._classify_slide_role(_VISUAL, False, "doc") == "visual", (
        "전제 실패: 사진형 슬라이드는 visual로 분류돼야 함")


# ==========================================================================
# SCENARIO 1 — HTML OFF: byte-level loss-zero + structural stays native
# ==========================================================================
def test_integration_mixed_deck_html_off():
    """Mixed deck, HTML off, Vertex on:
    (a) every generated Vertex image is embedded (unused == 0),
    (b) the structural slide stays editable native shapes (no full-bleed raster),
    (c) slide count == cover + 3 sections.
    """
    _assert_deck_classification()

    fake = _FakeVertexClient()
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "0",        # HTML off → Vertex embedded as bytes
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex pre-gen on
        "AE_PREFER_EDITABLE_DIAGRAM": "0",   # deterministic: no gateway structuring
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {"title": "통합 검증 덱", "slides": [dict(_STRUCTURAL), dict(_CONTENT), dict(_VISUAL)]}
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: None), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    pptx = result["absPath"]

    # (a) loss-zero — every generated Vertex image is embedded.
    generated, embedded, unused = _vertex_usage(fake, pptx)
    assert generated >= 2, (
        f"비구조형(content/visual) 슬라이드에 Vertex 이미지가 생성돼야 함 (generated={generated})")
    assert unused == 0, (
        f"손실-0 위반: 생성된 Vertex 이미지가 폐기됨 (generated={generated}, "
        f"embedded={embedded}, unused={unused})")

    # (c) slide count == cover + 3 sections.
    slides = _slides(pptx)
    assert len(slides) == 4, f"슬라이드 수 기대 4(표지+3) — 실제 {len(slides)}"

    # (b) structural slide (index 1) stays native: no full-bleed raster, bullets
    #     survive as editable shape text.
    struct_slide = slides[1]
    assert _fullbleed_pictures(struct_slide) == [], (
        "구조형 슬라이드에 풀블리드 래스터가 깔리면 안 됨(편집 가능 네이티브 도형이어야 함)")
    texts = " ".join(_shape_texts(struct_slide))
    hits = [b for b in _STRUCTURAL["bullets"] if b in texts]
    assert len(hits) >= 3, (
        f"구조형 라벨이 편집 가능 도형 텍스트로 남아야 함 — hits={hits}")


# ==========================================================================
# SCENARIO 2 — HTML ON: high-density HTML primary + Vertex coexistence
# ==========================================================================
def test_integration_mixed_deck_html_on_coexistence():
    """Mixed deck, HTML on, Vertex on:
    - Vertex generation is NOT suppressed by HTML (coexistence — fix to gate A),
    - high-density content slides take the full-bleed HTML slideBackground path,
    - the composited Vertex hero bytes surface in ppt/media (not discarded),
    - slide count == cover + 3 sections.
    """
    _assert_deck_classification()

    fake = _FakeVertexClient()
    captured_heroes: list = []
    proj = tempfile.mkdtemp()
    env = {
        "AE_ENABLE_HTML_SLIDES": "1",        # HTML on (primary renderer)
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "1",       # Vertex must coexist, not be suppressed
        "AE_PPTX_TOC": "0",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_DISABLE_NATIVE_COVER": "1",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {"title": "통합 검증 덱", "slides": [dict(_STRUCTURAL), dict(_CONTENT), dict(_VISUAL)]}
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: fake), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", lambda *a, **k: object()), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_render_html_slide_to_png", _render_html_png_fake), \
            patch.object(srv, "_generate_html_slide_for_section",
                         _make_compositing_section_fake(proj, captured_heroes)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx generation failed: {result}"
    pptx = result["absPath"]

    # Coexistence: HTML on did NOT suppress Vertex generation (fixed gate A).
    assert fake.calls >= 1, (
        "HTML 활성 시에도 Vertex 생성이 억제되면 안 됨(공존 게이트) — "
        f"vertex.generate 호출 {fake.calls}회")

    # The Vertex hero was fed into the HTML section renderer for non-structural
    # slides (proof the image is consumed by HTML compositing, not discarded).
    non_empty_heroes = [h for h in captured_heroes if h]
    assert non_empty_heroes, (
        "비구조형 슬라이드의 Vertex 히어로 이미지가 HTML 섹션 렌더러로 전달돼야 함(합성)")

    # Loss-zero under HTML: every generated Vertex image surfaces in media — the
    # compositing fake bakes hero bytes into the slideBackground PNG.
    generated, embedded, unused = _vertex_usage(fake, pptx)
    assert unused == 0, (
        f"손실-0 위반(HTML 경로): 생성된 Vertex 이미지가 최종 덱에 없음 "
        f"(generated={generated}, embedded={embedded}, unused={unused})")

    # High-density HTML path: content slides carry a (0,0) full-bleed background.
    slides = _slides(pptx)
    assert len(slides) == 4, f"슬라이드 수 기대 4(표지+3) — 실제 {len(slides)}"
    content_fb = _fullbleed_pictures(slides[2])  # index 2 == content slide
    assert content_fb, (
        "고밀도 콘텐츠 슬라이드는 HTML 풀블리드 슬라이드배경을 가져야 함")


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
