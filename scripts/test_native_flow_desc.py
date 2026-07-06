"""흐름도 단계 설명 텍스트 렌더링 검증."""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "ai_engine"))
from pptx import Presentation
from pptx.util import Inches
from native_diagram_pptx import build_native_diagram


def _texts(slide):
    out = []
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            out.append(sh.text_frame.text.strip())
    return out


def main():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # 설명 포함 가로 흐름(4단계)
    s1 = prs.slides.add_slide(prs.slide_layouts[6])
    c1 = ("요구사항 분석: 이해관계자 식별 및 문서화 -> "
          "설계: 시스템 아키텍처 및 DB 스키마 -> "
          "개발: 모듈 구현과 단위테스트 -> "
          "배포: CI/CD 파이프라인 운영")
    ok1 = build_native_diagram(s1, "flow", c1, palette=["#4472C4", "#ED7D31", "#A5A5A5"])
    t1 = " ".join(_texts(s1))
    print(f"[가로 4단계] ok={ok1} shapes={len(s1.shapes)}")
    assert ok1
    assert "이해관계자 식별" in t1, "설명 텍스트 누락(가로)"
    assert "CI/CD" in t1

    # 설명 포함 세로 흐름(6단계)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    c2 = "\n".join([
        "기획: 목표와 범위 정의",
        "분석: 요구사항 수집",
        "설계: 아키텍처 결정",
        "구현: 코드 작성",
        "테스트: 품질 검증",
        "배포: 운영 이행",
    ])
    ok2 = build_native_diagram(s2, "flow", c2, palette=["#4472C4", "#ED7D31"])
    t2 = " ".join(_texts(s2))
    print(f"[세로 6단계] ok={ok2} shapes={len(s2.shapes)}")
    assert ok2
    assert "요구사항 수집" in t2, "설명 텍스트 누락(세로)"
    assert "운영 이행" in t2

    # 설명 없는 흐름(라벨만) — 폴백 동작
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    ok3 = build_native_diagram(s3, "flow", "A -> B -> C", palette=["#4472C4"])
    print(f"[라벨만] ok={ok3} shapes={len(s3.shapes)}")
    assert ok3

    prs.save("/tmp/test_flow_desc.pptx")
    print("FLOW DESC PASS")


if __name__ == "__main__":
    main()
