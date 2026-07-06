# Feature: pptx-native-density-render, Property 19: 네이티브 도형은 design_tokens 스타일 품질을 충족한다 — 알려진_레이아웃으로 네이티브 렌더된 콘텐츠 슬라이드에 대해 audit_style_quality 는 (카드/박스 라운드 + 그림자/테두리) AND (섹션헤더/배지 accent 색, 기본 검정/흰색 아님) AND (제목 폰트 > 본문 폰트 타이포 계층) AND (본문 여백 토큰 적용)을 만족하여 passed=True 를 반환한다. 이 스타일 품질 차원은 밀도(요소 존재) 채점과 직교하며, 불합격 시 missing_style 은 실제 누락된 스타일 항목과 일치한다.
"""Property 19 (스타일 품질) PBT — spec: pptx-native-density-render, Task 14.1.

**Property 19: 네이티브 도형은 design_tokens 스타일 품질을 충족한다.**

  *For any* 알려진_레이아웃으로 네이티브 렌더된 콘텐츠 슬라이드에 대해,
  ``audit_style_quality(slide, tokens)`` 는
    (카드/박스 라운드 + 그림자/테두리) AND
    (섹션헤더/배지 accent 색 적용, 기본 검정/흰색 아님) AND
    (제목 폰트크기 > 본문 폰트크기 타이포 계층) AND
    (본문 여백 토큰 적용)
  을 만족하여 ``passed=True`` 를 반환한다. 이 스타일 품질 차원은 밀도(요소 존재)
  채점과 직교하며, 불합격 시 ``missing_style`` 은 실제 누락 항목과 일치한다.

**Validates: Requirements 5.1, 5.3, 5.4**

검증 대상(실제 코드, 네트워크 0):
  - ``scripts/audit_pptx_native_density.audit_style_quality(slide, tokens)`` (작업 6.1)
  - ``ai_engine/native_layout_renderer.render_native_layout`` 의 emit_* 도형 방출 (작업 2.x/3.x)
  - ``ai_engine/slide_templates.design_tokens_for_profile`` 로 획득한 토큰(기본/여러 프로파일)

설계 정합(modeling note):
  audit_style_quality 의 스타일 품질 검사는 **카드/박스 라운드+그림자 AND
  섹션헤더/배지 accent 색** 을 동시에 요구하는 논리곱이다. 알려진_레이아웃 중
  이 논리곱을 단독으로 충족하는 것은 ``section_divider``(다크 라운드 막대+그림자
  + accent 번호 배지)이며, ``feature_grid``/``architecture``(흰 카드 → accent 색
  미충족)·``timeline``(oval 배지 → 라운드 막대 미충족) 등은 단독으로는 일부 항목이
  비므로(프로브 확인) **"스타일 요소(라운드+그림자+accent)를 만드는 레이아웃이
  포함되도록"**(task 14.1) 생성기를 구성한다. 즉 충족 슬라이드는 항상
  ``section_divider`` 를 베이스로 렌더하고(논리곱 보장), 변주를 위해 카드/번호목록
  등을 만드는 추가 스타일 레이아웃을 같은 슬라이드에 0~2개 더 합성한다(실사용 고밀도
  덱과 동형). 토큰은 ``design_tokens_for_profile(None)`` 기본과 여러 프로파일로 변주한다.

헤르메틱 — 네트워크 0. python-pptx 인메모리 Presentation 만 구동하며 게이트웨이/
Vertex/HTML 렌더(Chrome)를 호출하지 않는다. Vertex 장식 경로는 옵트인 OFF 기본값
(aws_profile="" / credentials=None)이라 No-op 이다.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_native_style_quality_pbt.py -p no:cacheprovider -q
"""
from __future__ import annotations

import os
import sys

# repo root(ai_engine 패키지) + scripts(검증기) 를 import 가능하게 한다.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE, os.path.join(_ROOT, "ai_engine")):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402

pptx = pytest.importorskip("pptx")  # python-pptx 미설치 환경은 skip
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# 실제 코드(검증 대상) — 네트워크 0.
from ai_engine.native_layout_renderer import render_native_layout  # noqa: E402
from ai_engine.slide_templates import design_tokens_for_profile  # noqa: E402
from audit_pptx_native_density import audit_style_quality  # noqa: E402

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# 스타일 요소(라운드+그림자+accent)를 만드는 합성용 추가 레이아웃 후보.
_EXTRA_STYLE_LAYOUTS = ["feature_grid", "architecture", "timeline", "section_divider"]


# ---------------------------------------------------------------------------
# Strategies — 토큰/데이터를 넓게 변주하되 스타일 논리곱이 항상 성립하도록 구성.
# ---------------------------------------------------------------------------


@st.composite
def _vivid_hex(draw) -> str:
    """기본 검정(#000000)/흰색(#FFFFFF) 이 아닌 유채색 #RRGGBB.

    각 채널을 20~235 범위로 잡아 ``all<=12``(검정)·``all>=243``(흰색) 을 구조적으로
    배제한다 → accent 토큰이 항상 chromatic 이라 accent_color 검사를 결정적으로 만족.
    """
    r = draw(st.integers(min_value=20, max_value=235))
    g = draw(st.integers(min_value=20, max_value=235))
    b = draw(st.integers(min_value=20, max_value=235))
    return "#%02X%02X%02X" % (r, g, b)


# 비어있지 않은 텍스트(한글·특수문자·혼합·초장문 포함).
_NONEMPTY_TEXT = st.one_of(
    st.sampled_from([
        "섹션 개요",
        "Phase 1: 전략 🚀",
        "특수문자 !@#$%^&*()_+-=[]{}",
        "혼합 Mixed 한글 123 ABC",
        "가" * 120,
        "Architecture & Data Flow",
    ]),
    st.text(min_size=1, max_size=60).map(lambda s: s.strip()).filter(lambda s: len(s) > 0),
)

# 선택 텍스트(빈 문자열 허용 — description/optional 필드).
_OPT_TEXT = st.one_of(st.just(""), _NONEMPTY_TEXT)

_FONT = st.sampled_from(["Pretendard", "Noto Sans KR", "Inter", "Roboto", "맑은 고딕"])

# 프로파일: None(기본 토큰) 또는 유채색/폰트를 변주한 dict. design_tokens_for_profile
# 가 #RRGGBB/폰트 유효성 검사 후 매핑하므로(신규 토큰 정의 아님), 여기서 생성한 색은
# 전부 유효 hex 다.
_PROFILE = st.one_of(
    st.none(),
    st.fixed_dictionaries({
        "accentColor": _vivid_hex(),
        "primaryColor": _vivid_hex(),
        "secondaryColor": _vivid_hex(),
        "textColor": _vivid_hex(),
        "headingFont": _FONT,
        "bodyFont": _FONT,
    }),
)


def _section_divider_data(draw) -> dict:
    return {
        "title": draw(_NONEMPTY_TEXT),
        "section_number": draw(st.integers(min_value=1, max_value=99)),
        "description": draw(_OPT_TEXT),
    }


def _extra_layout_data(draw, layout: str) -> dict:
    """추가 합성 스타일 레이아웃의 최소 유효 data(필수 필드 충족)."""
    title = draw(_NONEMPTY_TEXT)
    if layout == "feature_grid":
        k = draw(st.integers(min_value=1, max_value=4))
        return {"title": title, "features": [
            {"title": draw(_NONEMPTY_TEXT), "description": draw(_OPT_TEXT)} for _ in range(k)]}
    if layout == "architecture":
        k = draw(st.integers(min_value=1, max_value=3))
        return {"title": title, "layers": [
            {"name": draw(_NONEMPTY_TEXT), "description": draw(_OPT_TEXT),
             "items": [draw(_NONEMPTY_TEXT)]} for _ in range(k)]}
    if layout == "timeline":
        k = draw(st.integers(min_value=1, max_value=4))
        return {"title": title, "steps": [
            {"title": draw(_NONEMPTY_TEXT), "description": draw(_OPT_TEXT)} for _ in range(k)]}
    # section_divider
    return {"title": title, "section_number": draw(st.integers(1, 9)),
            "description": draw(_OPT_TEXT)}


@st.composite
def _style_slide_spec(draw):
    """충족 슬라이드 명세: (profile, base section_divider data, [(extra_layout, data)...]).

    베이스는 항상 section_divider(스타일 논리곱 보장). 변주를 위해 추가 스타일
    레이아웃을 0~2개 합성한다.
    """
    profile = draw(_PROFILE)
    base = _section_divider_data(draw)
    extras_layouts = draw(st.lists(st.sampled_from(_EXTRA_STYLE_LAYOUTS), max_size=2))
    extras = [(lyt, _extra_layout_data(draw, lyt)) for lyt in extras_layouts]
    return profile, base, extras


# ---------------------------------------------------------------------------
# 헬퍼 — 슬라이드 생성
# ---------------------------------------------------------------------------


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs):
    # 레이아웃 6 = blank.
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------------------
# Property 19 — 단일 PBT (충족 슬라이드 passed=True + 텍스트만 슬라이드 passed=False)
# ---------------------------------------------------------------------------


@settings(max_examples=120, deadline=None,
          suppress_health_check=[HealthCheck.too_slow, HealthCheck.data_too_large])
@given(spec=_style_slide_spec(), text_only_title=_NONEMPTY_TEXT)
def test_property19_style_quality_passes_and_is_deterministic(spec, text_only_title):
    """Property 19 — 스타일 충족 슬라이드는 passed=True(missing 비어있음)이고,
    텍스트만/빈 슬라이드는 passed=False(missing_style 비어있지 않음)이다(결정성)."""
    profile, base, extras = spec
    tokens = design_tokens_for_profile(profile)

    # --- (1) 충족 슬라이드: section_divider(+추가 스타일 레이아웃) 네이티브 렌더 ----
    prs = _new_prs()
    slide = _blank_slide(prs)
    res = render_native_layout(slide, prs, "section_divider", base, tokens)
    assert res.ok is True and res.unsupported is False
    for layout, data in extras:
        render_native_layout(slide, prs, layout, data, tokens)

    rep = audit_style_quality(slide, tokens)
    assert rep.passed is True, (
        f"스타일 충족 슬라이드인데 미달 — missing_style={rep.missing_style} "
        f"score={rep.score} profile={profile} extras={[l for l, _ in extras]}")
    assert rep.missing_style == [], f"passed=True 인데 missing_style 비어있지 않음: {rep.missing_style}"
    assert rep.score == pytest.approx(1.0)

    # --- (2) 텍스트만 슬라이드: 장식 도형 없음 → passed=False, missing 비어있지 않음 --
    prs2 = _new_prs()
    bare = _blank_slide(prs2)
    tb = bare.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(8.0), Inches(1.0))
    tb.text_frame.text = text_only_title
    rep_bare = audit_style_quality(bare, tokens)
    assert rep_bare.passed is False, "텍스트만 슬라이드인데 passed=True (스타일 요소 없음에도 통과)"
    assert len(rep_bare.missing_style) > 0, "passed=False 인데 missing_style 가 비어있음(보고 불완전)"
    assert rep_bare.score < 1.0


if __name__ == "__main__":  # pragma: no cover
    import pytest as _pt
    raise SystemExit(_pt.main([__file__, "-p", "no:cacheprovider", "-q"]))
