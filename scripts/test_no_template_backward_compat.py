"""무템플릿 하위 호환 단위 테스트 (Validates: Requirements 5.2).

요구사항 5.2: 생성 요청에 templateId가 전달되지 않았거나 "템플릿 없음"이 선택된
경우, PPTX_Generator는 무템플릿 생성 경로(빈 `Presentation()` 기반)를 호출하며,
그 산출물은 동일 입력에 대한 기존(baseline) 산출물과 슬라이드 수·레이아웃 매핑·
배경 파이프라인 단계가 동일해야 한다.

이 단위 테스트는 `_tool_generate_pptx`를 templatePath/templateId/styleProfile을
전혀 전달하지 않고 호출하여 다음을 검증한다:

1. 무템플릿 생성    — 유효 PPTX 산출(path/.pptx, sizeBytes>0, error 없음),
                      응답에 templateId 키가 없어 기존 응답 형태 보존(요구사항 6.7 반대)
2. 슬라이드 수 일치  — 표지 1 + 입력 슬라이드 수
3. 슬라이드 크기 보존 — slide_width == Inches(13.333), slide_height == Inches(7.5)
                      (16:9, EMU 6858000) — 무템플릿 baseline 16:9 리사이즈 경로 보존
4. 레이아웃 매핑     — 무템플릿 경로에서 입력 layout("title"/"content"/"two-column")이
                      LAYOUT_MAP {title:0, content:1, two-column:3} 인덱스로 매핑됨
                      (빈 Presentation의 기본 레이아웃 세트 기준)

이미지 단계를 타지 않도록 슬라이드에는 imagePrompt/imageFile/slideBackground를
넣지 않는다 → Bedrock 호출 없음, 순수 결정론적 단위 테스트.

실행:
  ai_engine/.venv/bin/python scripts/test_no_template_backward_compat.py
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile

# 레포 루트에서 ai_engine 패키지를 import 가능하게 한다.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

# python-pptx 부재 시 skip (요구사항 9.3 환경과 별개로, 테스트는 환경 의존이므로 skip).
try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:  # pragma: no cover - 환경 의존
    print("SKIP: python-pptx 미설치 — 무템플릿 하위 호환 테스트를 건너뜁니다.")
    sys.exit(0)

from ai_engine.server import _tool_generate_pptx, LAYOUT_MAP  # noqa: E402


# 입력 슬라이드: title / content / two-column 레이아웃을 모두 포함한다.
# imagePrompt/imageFile/slideBackground 없음 → 텍스트 전용, 이미지 단계 미진입.
_INPUT_SLIDES = [
    {"title": "타이틀 레이아웃 슬라이드", "layout": "title", "bullets": ["요점 A", "요점 B"]},
    {"title": "콘텐츠 레이아웃 슬라이드", "layout": "content", "bullets": ["내용 1", "내용 2"]},
    {"title": "투-칼럼 레이아웃 슬라이드", "layout": "two-column", "bullets": ["좌측", "우측"]},
]


def _generate_no_template(tmp: str) -> tuple[dict, str]:
    """templatePath/templateId/styleProfile을 전혀 넣지 않고 _tool_generate_pptx 호출.

    Returns:
        (parsed_response_json, absolute_pptx_path)
    """
    tool_input = {
        "title": "무템플릿 하위 호환 덱",
        "slides": _INPUT_SLIDES,
        # templatePath / templateId / styleProfile 키를 의도적으로 전혀 넣지 않음
    }
    raw = asyncio.run(_tool_generate_pptx(tool_input, project_path=tmp))
    parsed = json.loads(raw)
    abs_path = os.path.join(tmp, parsed.get("path", ""))
    return parsed, abs_path


def test_no_template_generates_valid_pptx():
    """1. 무템플릿 입력 → 유효 PPTX 생성, 기존 응답 형태 보존."""
    with tempfile.TemporaryDirectory() as tmp:
        parsed, abs_path = _generate_no_template(tmp)

        assert "error" not in parsed, f"예상치 못한 에러: {parsed}"

        rel = parsed.get("path", "")
        assert rel.startswith(".generated/") and rel.endswith(".pptx"), (
            f"예상치 못한 path: {rel}"
        )
        assert os.path.isfile(abs_path), f"PPTX 미저장: {abs_path}"

        # sizeBytes > 0 (요구사항 6.7 baseline)
        assert parsed.get("sizeBytes", 0) > 0, f"sizeBytes>0 아님: {parsed}"
        assert os.path.getsize(abs_path) > 0, "PPTX 파일이 비어 있음"

        # 무템플릿 경로는 응답에 templateId를 포함하지 않아야 한다(하위 호환, 요구사항 5.2).
        assert "templateId" not in parsed, (
            f"무템플릿 응답에 templateId가 포함됨(기존 형태 미보존): {parsed}"
        )


def test_no_template_slide_count_matches():
    """2. 슬라이드 수 = 표지 1 + 입력 슬라이드 수."""
    with tempfile.TemporaryDirectory() as tmp:
        parsed, abs_path = _generate_no_template(tmp)

        expected = len(_INPUT_SLIDES) + 1  # +1 표지
        assert parsed.get("slideCount") == expected, (
            f"응답 slideCount 불일치: got {parsed.get('slideCount')}, expected {expected}"
        )

        prs = Presentation(abs_path)
        assert len(prs.slides) == expected, (
            f"파싱된 슬라이드 수 {len(prs.slides)} != 기대 {expected}"
        )


def test_no_template_slide_size_preserved():
    """3. 슬라이드 크기 16:9 보존 — Inches(13.333) × Inches(7.5) (EMU 6858000)."""
    with tempfile.TemporaryDirectory() as tmp:
        _parsed, abs_path = _generate_no_template(tmp)

        prs = Presentation(abs_path)
        assert prs.slide_width == Inches(13.333), (
            f"slide_width EMU 불일치: got {prs.slide_width}, "
            f"expected {int(Inches(13.333))}"
        )
        assert prs.slide_height == Inches(7.5), (
            f"slide_height EMU 불일치: got {prs.slide_height}, "
            f"expected {int(Inches(7.5))}"
        )
        # 7.5 inch == 6858000 EMU (명시 검증)
        assert int(prs.slide_height) == 6858000, (
            f"slide_height != 6858000 EMU: {int(prs.slide_height)}"
        )


def test_no_template_layout_map_preserved():
    """4. 무템플릿 경로에서 LAYOUT_MAP {title:0, content:1, two-column:3} 매핑 보존.

    빈 Presentation의 기본 레이아웃 세트 기준:
      - 표지(슬라이드 0): "title" → LAYOUT_MAP["title"]=0
      - 입력 슬라이드 i: 해당 layout → LAYOUT_MAP[layout]

    저장된 .pptx를 다시 열어 각 슬라이드의 slide_layout이 LAYOUT_MAP가 가리키는
    인덱스의 레이아웃과 동일한지 검증한다.
    """
    # 전제 확인: LAYOUT_MAP가 기대 매핑과 일치(구현 상수 회귀 방지).
    assert LAYOUT_MAP == {"title": 0, "content": 1, "two-column": 3}, (
        f"LAYOUT_MAP가 기대 매핑과 다름: {LAYOUT_MAP}"
    )

    with tempfile.TemporaryDirectory() as tmp:
        _parsed, abs_path = _generate_no_template(tmp)

        prs = Presentation(abs_path)
        layouts = list(prs.slide_layouts)

        def _layout_index(slide):
            # slide.slide_layout이 prs.slide_layouts 컬렉션에서 차지하는 인덱스.
            return layouts.index(slide.slide_layout)

        slides = list(prs.slides)

        # 표지(슬라이드 0)는 "title" 레이아웃으로 해석 → 인덱스 0.
        assert _layout_index(slides[0]) == LAYOUT_MAP["title"], (
            f"표지 레이아웃 인덱스 {_layout_index(slides[0])} != "
            f"LAYOUT_MAP['title']={LAYOUT_MAP['title']}"
        )

        # 이후 입력 슬라이드들의 레이아웃 인덱스가 LAYOUT_MAP과 일치.
        for i, sd in enumerate(_INPUT_SLIDES):
            slide = slides[i + 1]  # +1: 표지 다음부터
            expected_idx = LAYOUT_MAP[sd["layout"]]
            actual_idx = _layout_index(slide)
            assert actual_idx == expected_idx, (
                f"슬라이드 {i + 1} (layout={sd['layout']}) 레이아웃 인덱스 불일치: "
                f"got {actual_idx}, expected {expected_idx}"
            )


def main():
    print("=== 무템플릿 하위 호환 단위 테스트 (요구사항 5.2) ===")
    test_no_template_generates_valid_pptx()
    print("  1. 무템플릿 → 유효 PPTX + 응답 형태 보존        OK")
    test_no_template_slide_count_matches()
    print("  2. 슬라이드 수 = 표지 1 + 입력 수                OK")
    test_no_template_slide_size_preserved()
    print("  3. 슬라이드 크기 13.333×7.5in (6858000 EMU) 보존 OK")
    test_no_template_layout_map_preserved()
    print("  4. LAYOUT_MAP {title:0,content:1,two-column:3} 보존 OK")
    print("모든 무템플릿 하위 호환 케이스 통과.")


if __name__ == "__main__":
    main()
