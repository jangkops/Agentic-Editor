"""Hermetic property-based tests — spec: pptx-ultra-quality-hybrid-render, Task 7.5/7.6/7.7.

Feature: pptx-ultra-quality-hybrid-render

These tests exercise the REAL ``ai_engine/server.py:_tool_generate_pptx`` hybrid
decision seam end to end, fully hermetically (no network):

  - Bedrock gateway (`_get_gw`) is never reached (HTML off + image tool disabled).
  - Vertex client (`vertex_image_module.get_vertex_image_client`) is replaced with a
    spy/disabled stub — its ``generate`` records calls and returns local PNG bytes.
  - The HTML→PNG render path is disabled (``AE_DISABLE_HTML_SLIDES=1`` +
    ``_find_local_chrome`` → "" + ``_call_bridge`` → None) so no Chrome/bridge runs.
  - ``_tool_generate_image`` is mocked to never touch Bedrock.

Mocking mirrors ``scripts/test_pptx_quality_vertex_images_integration.py``.

Properties implemented (design.md):

  Property 13 (Task 7.5) — Vertex ``generate`` 호출 계약 (Validates R3.1/3.5):
      풀블리드 대상 슬라이드(caller 미지정, role∈{cover,section,visual}, vertex_enabled)에서
      ``VertexImageClient.generate`` 가 ``aspect_ratio=="16:9"`` ∧
      ``model_class=="image_generation_high_quality"`` 로 정확히 1회 호출됨.

  Property 21 (Task 7.6) — Flag off 결정성 (Validates R6.4/1.7):
      ``AE_HYBRID_RENDER`` 비활성에서 동일 입력 반복 렌더가 구조적으로 동등한(슬라이드 수·
      도형·텍스트·이미지 배치 동일) 결정론적 산출물을 내고, 하이브리드 플랜 선택기
      ``_select_hybrid_render_plan`` 이 호출되지 않음(스파이).

  Property 19 (Task 7.7) — 산출물 겹침·편집성 감사 통과 (Validates R5.1):
      하이브리드 렌더(``AE_HYBRID_RENDER=1``, Vertex 비활성)로 생성한 .pptx 에 대해
      ``scripts/audit_pptx_overlap.py`` 의 ``audit(path)`` 가 판정한 "텍스트·이미지 겹침
      슬라이드" 개수 == 0 ∧ "편집 불가(래스터) 의심 슬라이드" 개수 == 0.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_hybrid_render_wiring_pbt.py -p no:cacheprovider -q
"""
from __future__ import annotations

import io
import os
import sys
import json
import base64
import asyncio
import tempfile
import importlib.util
from unittest.mock import patch

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402
from hypothesis import given, settings, strategies as st, assume, HealthCheck  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402

# audit_pptx_overlap.audit — import by file path (script has no package).
_AUDIT_PATH = os.path.join(os.path.dirname(__file__), "audit_pptx_overlap.py")
_spec = importlib.util.spec_from_file_location("_audit_overlap_mod", _AUDIT_PATH)
_audit_mod = importlib.util.module_from_spec(_spec)
_spec.loader.exec_module(_audit_mod)
audit_overlap = _audit_mod.audit


# --------------------------------------------------------------------------
# Hermetic fakes
# --------------------------------------------------------------------------
def _make_png(tag: int) -> bytes:
    img = Image.new("RGB", (48 + (tag % 9), 36 + (tag % 5)),
                    (tag % 256, (tag * 7) % 256, (tag * 31) % 256))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


class _SpyVertexClient:
    """VertexImageClient stand-in. Records every generate() call's kwargs."""

    def __init__(self, enabled: bool = True) -> None:
        self.enabled = enabled
        self.calls: list[dict] = []

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None,
                       negative_prompt=None, timeout=None, **_kw):
        self.calls.append({
            "prompt": prompt,
            "model_class": model_class,
            "aspect_ratio": aspect_ratio,
            "negative_prompt": negative_prompt,
            "timeout": timeout,
        })
        raw = _make_png(1000 + len(self.calls))
        return {"images": [base64.b64encode(raw).decode("ascii")]}


async def _img_gen_disabled(*_a, **_k):
    return json.dumps({"error": "test-disabled"})


def _base_env(proj: str) -> dict:
    """Env that guarantees a fully hermetic, HTML-off render."""
    return {
        "AE_DISABLE_HTML_SLIDES": "1",   # force HTML full-bleed bake OFF (no Chrome)
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
        # AE_ENABLE_VERTEX_BG unset → the direct-call cover/body Vertex bg path is OFF.
        # AE_PREFER_VERTEX_IMAGE default "1" → per-slide Vertex pre-gen loop active.
    }


def _run_pptx(deck: dict, env: dict, vertex_client, extra_patches=None) -> dict:
    """Drive the real _tool_generate_pptx hermetically; return parsed JSON result."""
    proj = env["AE_GENERATED_ROOT"]
    patches = [
        patch.dict(os.environ, env, clear=False),
        patch.object(vim, "get_vertex_image_client", lambda **_k: vertex_client),
        patch.object(srv, "_call_bridge", lambda *a, **k: None),
        patch.object(srv, "_find_local_chrome", lambda: ""),
        patch.object(srv, "_tool_generate_image", _img_gen_disabled),
    ]
    for p in (extra_patches or []):
        patches.append(p)
    import contextlib
    with contextlib.ExitStack() as stack:
        for p in patches:
            stack.enter_context(p)
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    return json.loads(raw)


# --------------------------------------------------------------------------
# Structural signature helpers (Property 21)
# --------------------------------------------------------------------------
def _slide_signature(slide) -> tuple:
    """Structure-only fingerprint: (shape_count, sorted texts, sorted pic geometries)."""
    texts = []
    pics = []
    n = 0
    for sh in slide.shapes:
        n += 1
        try:
            st_ = sh.shape_type
        except Exception:
            st_ = None
        if st_ == MSO_SHAPE_TYPE.PICTURE:
            try:
                pics.append((int(sh.left or 0), int(sh.top or 0),
                             int(sh.width or 0), int(sh.height or 0)))
            except Exception:
                pics.append((None, None, None, None))
        else:
            try:
                if getattr(sh, "has_text_frame", False):
                    t = (sh.text_frame.text or "").strip()
                    if t:
                        texts.append(t)
            except Exception:
                continue
    return (n, tuple(sorted(texts)), tuple(sorted(pics)))


def _deck_signature(pptx_abs: str) -> tuple:
    prs = Presentation(pptx_abs)
    return tuple(_slide_signature(s) for s in prs.slides)


# --------------------------------------------------------------------------
# Slide generators
# --------------------------------------------------------------------------
_SAFE_WORDS = ["회사", "비전", "성장", "혁신", "미래", "고객", "가치", "전략",
               "브랜드", "신뢰", "품질", "도전", "성과", "협력", "지속가능"]

_topic_st = st.lists(st.sampled_from(_SAFE_WORDS), min_size=1, max_size=3).map(
    lambda ws: " ".join(ws))


@st.composite
def _visual_slide(draw):
    """A slide that deterministically classifies as role == 'visual'."""
    topic = draw(_topic_st)
    subtitle = draw(st.sampled_from(["신뢰를 최우선으로", "미래를 향해", "함께 성장",
                                      "새로운 도약", "고객 중심"]))
    prompt_lead = draw(st.sampled_from([
        "a high quality professional photograph of",
        "cinematic professional photograph of",
        "a polished editorial photograph of",
    ]))
    return {
        "title": topic,
        "bullets": [subtitle],
        "imagePrompt": f"{prompt_lead} {topic}, natural light, wide angle, corporate",
    }


@st.composite
def _content_slide(draw):
    """A slide that deterministically classifies as role == 'content' (dense bullets)."""
    title = draw(_topic_st)
    n = draw(st.integers(min_value=4, max_value=6))
    pool = ["빠른 처리", "안정적 운영", "유연한 확장", "강력한 보안",
            "비용 절감", "쉬운 사용성", "높은 가용성", "실시간 분석"]
    return {"title": f"{title} 요약", "bullets": pool[:n]}


@st.composite
def _structural_slide(draw):
    """A slide that deterministically classifies as role == 'structural' (flow)."""
    title = draw(st.sampled_from(["업무 처리 프로세스", "승인 흐름", "배포 파이프라인"]))
    return {"title": title, "bullets": ["접수", "검토", "승인", "완료"]}


# ==========================================================================
# Property 13 (Task 7.5) — Vertex generate 호출 계약
# **Validates: Requirements 3.1, 3.5**
# ==========================================================================
@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.too_slow, HealthCheck.function_scoped_fixture])
@given(vslide=_visual_slide(), deck_title=_topic_st)
def test_property13_vertex_generate_contract(vslide, deck_title):
    """풀블리드 대상(role=visual, caller 미지정, vertex_enabled) 슬라이드에서
    VertexImageClient.generate 가 aspect_ratio=='16:9' ∧
    model_class=='image_generation_high_quality' 로 정확히 1회 호출된다."""
    # 전제: 이 슬라이드는 반드시 풀블리드 대상(role∈{cover,section,visual})으로 분류.
    assume(srv._classify_slide_role(vslide, False, deck_title) in ("cover", "section", "visual"))

    spy = _SpyVertexClient(enabled=True)
    proj = tempfile.mkdtemp()
    env = _base_env(proj)
    env["AE_HYBRID_RENDER"] = "1"
    env["AE_PREFER_VERTEX_IMAGE"] = "1"

    # 단일 풀블리드 대상 body 슬라이드 → 사전생성 루프가 정확히 1회 generate 호출.
    deck = {"title": deck_title or "덱", "slides": [dict(vslide)]}
    result = _run_pptx(deck, env, spy)
    assert "absPath" in result, f"pptx 생성 실패: {result}"

    # 정확히 1회 호출.
    assert len(spy.calls) == 1, (
        f"풀블리드 대상 슬라이드는 generate 를 정확히 1회 호출해야 함 — 실제 {len(spy.calls)}회")
    call = spy.calls[0]
    assert call["aspect_ratio"] == "16:9", (
        f"aspect_ratio 는 '16:9' 여야 함 — 실제 {call['aspect_ratio']!r}")
    assert call["model_class"] == "image_generation_high_quality", (
        f"model_class 는 'image_generation_high_quality' 여야 함 — 실제 {call['model_class']!r}")


# ==========================================================================
# Property 21 (Task 7.6) — Flag off 결정성 (기존 동작 보존)
# **Validates: Requirements 6.4, 1.7**
# ==========================================================================
@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.too_slow, HealthCheck.function_scoped_fixture])
@given(flag=st.just("0"),
       slides=st.lists(st.one_of(_content_slide(), _structural_slide()),
                       min_size=1, max_size=2),
       deck_title=_topic_st)
def test_property21_flag_off_determinism(flag, slides, deck_title):
    """AE_HYBRID_RENDER 킬스위치("0")에서 (a) 하이브리드 플랜 선택기가 호출되지 않고,
    (b) 동일 입력 반복 렌더가 구조적으로 동등한 결정론적 산출물을 낸다.

    NOTE: 기본 ON(A안)으로 계약이 반전됨 — 미설정/"" 은 이제 하이브리드 ON 이므로
    legacy/off 경로는 명시적 킬스위치 AE_HYBRID_RENDER="0" 으로만 도달한다."""
    proj = tempfile.mkdtemp()
    env = _base_env(proj)
    # 킬스위치 명시: legacy(off) 경로는 "0" 으로만 진입한다(미설정/"" 는 이제 기본 ON).
    env["AE_HYBRID_RENDER"] = flag

    # Vertex 비활성 → 결정론(이미지 사전생성 없음). 스파이는 enabled=False.
    disabled_vertex = _SpyVertexClient(enabled=False)

    real_hybrid_plan = srv._select_hybrid_render_plan
    hybrid_calls = {"n": 0}

    def _spy_hybrid_plan(*a, **k):
        hybrid_calls["n"] += 1
        return real_hybrid_plan(*a, **k)

    deck = {"title": deck_title or "덱", "slides": [dict(s) for s in slides]}
    spy_patch = patch.object(srv, "_select_hybrid_render_plan", _spy_hybrid_plan)

    r1 = _run_pptx(deck, dict(env), disabled_vertex, extra_patches=[spy_patch])
    assert "absPath" in r1, f"1차 pptx 생성 실패: {r1}"
    sig1 = _deck_signature(r1["absPath"])

    # 하이브리드 플랜 선택기는 flag off 에서 절대 호출되지 않아야 한다.
    assert hybrid_calls["n"] == 0, (
        f"flag off 에서 _select_hybrid_render_plan 이 호출됨 ({hybrid_calls['n']}회) — "
        f"하이브리드 분기가 no-op 이 아님")

    # 2차 렌더 — 동일 입력.
    disabled_vertex2 = _SpyVertexClient(enabled=False)
    spy_patch2 = patch.object(srv, "_select_hybrid_render_plan", _spy_hybrid_plan)
    r2 = _run_pptx(deck, dict(env), disabled_vertex2, extra_patches=[spy_patch2])
    assert "absPath" in r2, f"2차 pptx 생성 실패: {r2}"
    sig2 = _deck_signature(r2["absPath"])

    assert r1["slideCount"] == r2["slideCount"], (
        f"슬라이드 수 불일치: {r1['slideCount']} vs {r2['slideCount']}")
    assert sig1 == sig2, (
        "flag off 반복 렌더가 구조적으로 동등하지 않음(슬라이드 수·도형·텍스트·이미지 배치 불일치)")


# ==========================================================================
# Property 19 (Task 7.7) — 산출물 겹침·편집성 감사 통과
# **Validates: Requirements 5.1**
# ==========================================================================
@settings(max_examples=100, deadline=None,
          suppress_health_check=[HealthCheck.too_slow, HealthCheck.function_scoped_fixture])
@given(slides=st.lists(st.one_of(_content_slide(), _structural_slide(), _visual_slide()),
                       min_size=1, max_size=3),
       deck_title=_topic_st)
def test_property19_overlap_and_editability_audit(slides, deck_title):
    """하이브리드 렌더(AE_HYBRID_RENDER=1, Vertex 비활성)로 생성한 .pptx 에 대해
    audit_pptx_overlap.audit(path) 의 겹침 슬라이드 목록 개수 == 0 ∧
    편집 불가(래스터) 의심 슬라이드 목록 개수 == 0."""
    proj = tempfile.mkdtemp()
    env = _base_env(proj)
    env["AE_HYBRID_RENDER"] = "1"

    # Vertex 비활성 — 헤르메틱(자격증명 부재) 하이브리드 경로. 편집 가능 네이티브 산출.
    disabled_vertex = _SpyVertexClient(enabled=False)

    deck = {"title": deck_title or "덱", "slides": [dict(s) for s in slides]}
    result = _run_pptx(deck, env, disabled_vertex)
    assert "absPath" in result, f"pptx 생성 실패: {result}"

    grand = audit_overlap(result["absPath"])
    assert isinstance(grand, dict), "audit(path) 는 판정 결과 dict 를 반환해야 함"

    overlap_ct = len(set(grand.get("overlap_slides", [])))
    baked_ct = len(set(grand.get("text_baked_slides", [])))
    assert overlap_ct == 0, (
        f"텍스트·이미지 겹침 슬라이드 개수 != 0 (개수={overlap_ct}, "
        f"슬라이드={sorted(set(grand.get('overlap_slides', [])))})")
    assert baked_ct == 0, (
        f"편집 불가(래스터) 의심 슬라이드 개수 != 0 (개수={baked_ct}, "
        f"슬라이드={sorted(set(grand.get('text_baked_slides', [])))})")


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
