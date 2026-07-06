"""동적 레이아웃 폴백 체인 단위 테스트 (Task 7.4).

server.py의 동적 레이아웃 매핑(설계 §구성요소 4)을 검증한다. 구현 코드는 수정하지 않고,
이미 구현된 다음 심볼만 사용한다.

  - _resolve_layout(prs, layout_name, used_template) -> slide_layout
  - _layout_has_content_placeholder(layout) -> bool
  - _layout_name_matches(name_norm, layout_name) -> bool
  - LAYOUT_MAP = {"title": 0, "content": 1, "two-column": 3}

검증 항목:
  1. 이름 매칭 성공 (요구사항 6.2)
     템플릿 레이아웃 이름이 title/content/two-column 의미로 정규화 매칭되면 그 레이아웃을 선택.
  2. 이름 매칭 실패 → 첫 콘텐츠 레이아웃 (요구사항 6.3)
     요청 layout에 대응하는 이름이 없으면 body placeholder를 가진 첫 콘텐츠 레이아웃으로 폴백.
  3. 콘텐츠 레이아웃 없음 → index 0 (요구사항 6.4)
     콘텐츠 레이아웃이 하나도 없으면 slide_layouts[0] 사용.
  4. 무템플릿 경로 LAYOUT_MAP {title:0, content:1, two-column:3} 사용 확인 (요구사항 5.2).

fixture 전략:
  - 이름 매칭(6.2) / 무템플릿 LAYOUT_MAP(5.2)은 python-pptx 기본 Presentation()의 표준
    레이아웃 세트(Title Slide / Title and Content / Two Content ...)로 실제 동작을 검증한다.
  - "첫 콘텐츠 레이아웃 폴백"(6.3) / "콘텐츠 레이아웃 없음"(6.4)은 표준 세트로 재현하기
    어려우므로, 레이아웃 목록을 stub하여 _resolve_layout의 폴백 체인을 격리 검증한다.

python-pptx import 불가 시 전체 스킵.

실행:
  ai_engine/.venv/bin/python scripts/test_dynamic_layout_mapping.py
"""
from __future__ import annotations

import os
import sys

# Make the ai_engine package importable from the repo root (기존 scripts 컨벤션).
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

# python-pptx 미설치 시 적절히 스킵 (요구사항 9.3 환경 보호).
try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
except Exception as e:  # pragma: no cover - 환경 의존
    print(f"SKIP: python-pptx import 불가 → {e}")
    sys.exit(0)

from ai_engine.server import (  # noqa: E402
    LAYOUT_MAP,
    _layout_has_content_placeholder,
    _layout_name_matches,
    _resolve_layout,
)


# --------------------------------------------------------------------------- #
# Stub 레이아웃 (폴백 체인 격리용)
#
# _resolve_layout / _layout_has_content_placeholder 가 실제로 접근하는 속성만 흉내낸다.
#   layout.name                                  (이름 매칭)
#   layout.placeholders -> [ph, ...]             (콘텐츠 placeholder 탐지)
#   ph.placeholder_format.idx / .type            (idx>=1 AND type in {BODY,OBJECT})
#   prs.slide_layouts -> list (indexable + iterable)
# --------------------------------------------------------------------------- #

class _FakePF:
    def __init__(self, idx, ptype):
        self.idx = idx
        self.type = ptype


class _FakePlaceholder:
    def __init__(self, idx, ptype):
        self.placeholder_format = _FakePF(idx, ptype)


class _FakeLayout:
    def __init__(self, name, placeholders=None):
        self.name = name
        self.placeholders = placeholders or []


class _FakePrs:
    def __init__(self, layouts):
        # 실제 python-pptx의 slide_layouts는 인덱싱과 반복을 모두 지원한다. list가 둘 다 만족.
        self.slide_layouts = layouts


def _title_only_layout(name):
    """본문 placeholder가 없는(제목만) 레이아웃 stub → 콘텐츠 레이아웃이 아니다."""
    return _FakeLayout(name, [_FakePlaceholder(0, PP_PLACEHOLDER.TITLE)])


def _content_layout(name):
    """idx>=1 본문(BODY) placeholder를 가진 콘텐츠 레이아웃 stub."""
    return _FakeLayout(
        name,
        [_FakePlaceholder(0, PP_PLACEHOLDER.TITLE),
         _FakePlaceholder(1, PP_PLACEHOLDER.BODY)],
    )


# --------------------------------------------------------------------------- #
# 1. 이름 매칭 성공 (요구사항 6.2) — 실제 기본 Presentation 사용
# --------------------------------------------------------------------------- #

def test_name_match_success_real_presentation():
    """기본 Presentation()의 표준 레이아웃 이름이 의미별로 매칭되어 그 레이아웃을 선택한다."""
    prs = Presentation()
    names = [l.name for l in prs.slide_layouts]
    # 사전 조건: 표준 세트가 기대한 이름을 포함하는지 확인 (환경 변화 시 조기 실패).
    assert "Title Slide" in names, names
    assert "Title and Content" in names, names
    assert "Two Content" in names, names

    by_name = {l.name: l for l in prs.slide_layouts}

    # title → "Title Slide" (표지류; "...and Content"는 제외되어야 함)
    got_title = _resolve_layout(prs, "title", used_template=True)
    assert got_title is by_name["Title Slide"], (
        f"title 매칭 실패: got name={got_title.name!r}"
    )

    # content → "Title and Content"
    got_content = _resolve_layout(prs, "content", used_template=True)
    assert got_content is by_name["Title and Content"], (
        f"content 매칭 실패: got name={got_content.name!r}"
    )

    # two-column → "Two Content"
    got_two = _resolve_layout(prs, "two-column", used_template=True)
    assert got_two is by_name["Two Content"], (
        f"two-column 매칭 실패: got name={got_two.name!r}"
    )


def test_layout_name_matches_unit():
    """_layout_name_matches 의 의미 판정 단위 검증 (요구사항 6.2)."""
    # title: 표지/제목류는 매칭, '제목 및 내용'류(title+content)는 제외
    assert _layout_name_matches("title slide", "title") is True
    assert _layout_name_matches("표지", "title") is True
    assert _layout_name_matches("cover", "title") is True
    assert _layout_name_matches("title and content", "title") is False  # content 포함 → 제외
    assert _layout_name_matches("제목 및 내용", "title") is False

    # content: 제목+내용 / content / 내용
    assert _layout_name_matches("title and content", "content") is True
    assert _layout_name_matches("제목 및 내용", "content") is True
    assert _layout_name_matches("blank", "content") is False

    # two-column: two content / comparison / 비교 / 2단
    assert _layout_name_matches("two content", "two-column") is True
    assert _layout_name_matches("comparison", "two-column") is True
    assert _layout_name_matches("비교", "two-column") is True
    assert _layout_name_matches("2단 레이아웃", "two-column") is True
    assert _layout_name_matches("title slide", "two-column") is False


# --------------------------------------------------------------------------- #
# 2. 이름 매칭 실패 → 첫 콘텐츠 레이아웃 (요구사항 6.3) — stub 격리
# --------------------------------------------------------------------------- #

def test_name_mismatch_falls_back_to_first_content_layout():
    """요청 layout에 대응하는 이름이 없으면 body placeholder를 가진 '첫' 콘텐츠 레이아웃 선택."""
    l0 = _title_only_layout("Alpha Cover")   # 콘텐츠 아님 + 이름 비매칭
    l1 = _content_layout("Beta Section")     # 첫 콘텐츠 레이아웃
    l2 = _content_layout("Gamma Section")    # 또 다른 콘텐츠 레이아웃
    prs = _FakePrs([l0, l1, l2])

    # "two-column" 의미에 매칭되는 이름이 하나도 없음 → 첫 콘텐츠 레이아웃(l1)로 폴백.
    got = _resolve_layout(prs, "two-column", used_template=True)
    assert got is l1, f"첫 콘텐츠 레이아웃 폴백 실패: got name={getattr(got, 'name', None)!r}"
    # 콘텐츠 placeholder가 없는 l0은 선택되지 않아야 한다.
    assert got is not l0


def test_layout_has_content_placeholder_unit():
    """_layout_has_content_placeholder 단위 검증 (요구사항 6.3)."""
    assert _layout_has_content_placeholder(_content_layout("c")) is True
    # 제목만 있는 레이아웃(idx 0 TITLE)은 콘텐츠가 아니다.
    assert _layout_has_content_placeholder(_title_only_layout("t")) is False
    # placeholder가 전혀 없는 레이아웃도 콘텐츠가 아니다.
    assert _layout_has_content_placeholder(_FakeLayout("empty", [])) is False


# --------------------------------------------------------------------------- #
# 3. 콘텐츠 레이아웃 없음 → index 0 (요구사항 6.4) — stub 격리
# --------------------------------------------------------------------------- #

def test_no_content_layout_falls_back_to_index_zero():
    """이름 매칭 실패 + 콘텐츠 레이아웃이 하나도 없으면 slide_layouts[0]을 사용한다."""
    l0 = _title_only_layout("Xray")   # 콘텐츠 아님
    l1 = _FakeLayout("Yankee", [])    # placeholder 없음
    prs = _FakePrs([l0, l1])

    got = _resolve_layout(prs, "content", used_template=True)
    assert got is l0, f"index 0 폴백 실패: got name={getattr(got, 'name', None)!r}"
    assert got is prs.slide_layouts[0]


# --------------------------------------------------------------------------- #
# 4. 무템플릿 경로 LAYOUT_MAP 사용 (요구사항 5.2) — 실제 기본 Presentation 사용
# --------------------------------------------------------------------------- #

def test_no_template_uses_hardcoded_layout_map():
    """used_template=False면 LAYOUT_MAP {title:0, content:1, two-column:3} 인덱스를 그대로 사용."""
    assert LAYOUT_MAP == {"title": 0, "content": 1, "two-column": 3}, LAYOUT_MAP
    prs = Presentation()

    assert _resolve_layout(prs, "title", used_template=False) is prs.slide_layouts[0]
    assert _resolve_layout(prs, "content", used_template=False) is prs.slide_layouts[1]
    assert _resolve_layout(prs, "two-column", used_template=False) is prs.slide_layouts[3]
    # 알 수 없는 layout 이름은 기본 content(인덱스 1)로 매핑.
    assert _resolve_layout(prs, "unknown-xyz", used_template=False) is prs.slide_layouts[1]


def main():
    print("=== Task 7.4: 동적 레이아웃 폴백 체인 단위 테스트 ===")
    cases = [
        ("이름 매칭 성공 (6.2, 실제 Presentation)", test_name_match_success_real_presentation),
        ("_layout_name_matches 단위 (6.2)", test_layout_name_matches_unit),
        ("이름 매칭 실패 → 첫 콘텐츠 레이아웃 (6.3, stub)", test_name_mismatch_falls_back_to_first_content_layout),
        ("_layout_has_content_placeholder 단위 (6.3)", test_layout_has_content_placeholder_unit),
        ("콘텐츠 레이아웃 없음 → index 0 (6.4, stub)", test_no_content_layout_falls_back_to_index_zero),
        ("무템플릿 LAYOUT_MAP 사용 (5.2, 실제 Presentation)", test_no_template_uses_hardcoded_layout_map),
    ]
    for label, fn in cases:
        fn()
        print(f"  {label:<48} OK")
    print(f"All {len(cases)} dynamic-layout cases passed.")


if __name__ == "__main__":
    main()
