"""Bug-condition exploration test — spec: pptx-overlay-collision-fix (bugfix), Task 1.

PROPERTY 1 — Bug Condition (오버레이/충돌 결함 A/B/C 재현).

This is an EXPLORATORY bug-condition test. It encodes the EXPECTED (post-fix)
behaviour — *after layout, every (textBox ∪ badge) pair must overlap less than
10% of the smaller box, a structural slide must NOT be a baked-text image
carrier, and a body box must overlap a full-bleed background by < 10%* (design
Property 1·2 / Fix Checking). On the UNFIXED code the placement/render seams in
``ai_engine/native_diagram_pptx.py`` (``build_native_cover`` / twocol ``_badge``)
and ``ai_engine/server.py`` (``_classify_slide_role`` + the ``_eff_bg`` full-bleed
composition) violate this, so these assertions FAIL — and that failure is the
proof the bug exists. After the fix (tasks 3.x) the SAME test is re-run (task 3.7)
and must PASS.

DO NOT "fix" this test or the code when it fails on unfixed code — the failure is
intended. It surfaces the counterexamples that prove defects A / B / C.

Scoped (deterministic) reproduction cases — design Examples coordinates:

  (A) defectA — 표지 제목↔부제 수직 겹침
      ``build_native_cover`` with a ``title_pt>=40`` title: title box bottom
      (top+height = 2.8 + 2.0 = 4.8) > subtitle top (title_y + 1.05 = 3.85) →
      vertical overlap ≈ 10.05 in² (~95% of the smaller box). The fixed-offset
      ``sub_y`` ignores the title box's real occupied height → defectA true.

  (C) defectC — 번호 배지↔라벨 박스 겹침
      twocol native cards (6 items): ``_badge(x0+3.4, …)`` is centred *inside* the
      label card (x0..x0+col_w) → badge rect fully contained → badge ∩ label ==
      100% of the badge area → defectC true.

  (B-1) defectB — 구조형이 구워진-텍스트 풀블리드 이미지로 렌더되지 않아야 함
      수정된 결정 seam ``_select_render_plan(role="structural", has_vertex_image=True,
      …, bg_has_baked_text=True)`` 를 구동 → ``primary=NATIVE_SHAPES`` ∧
      ``vertex_slot=backdrop``(손실-0) ∧ ``body_separated=True``. 또한
      ``_classify_slide_role(slide, is_cover=False, bg_has_baked_text=True)`` 가 시각형
      슬라이드를 visual→content 로 강등(본문 캐리어 차단). 수정 전(파라미터/키 부재)에는
      TypeError/KeyError 로 FAIL → fix 검증.

  (B-2) defectB — 풀블리드 배경 위 본문↔배경 분리
      수정된 순수 기하 ``layout_geometry.body_safe_area(slide, bg, has_baked_text=True,
      desired)`` 를 구동 → 반환 안전 영역 region' 에 대해
      ``overlap_area(region', bg) < 0.10*area(region')``. ``body_safe_area`` 가 없던 수정
      전엔 import 자체가 불가 → fix 검증. 배경이 슬라이드를 완전히 덮지 않는 현실적 입력으로
      구성해 본문이 배경 미점유 여백 띠로 분리되는 경로를 검증한다.

Hermetic — no network. We only drive the *native assembly* path
(``build_native_cover`` / ``build_native_diagram``) plus the pure decision
function ``_classify_slide_role`` and a locally-generated baked-text PNG. The
Bedrock gateway, the Vertex client and the HTML→PNG renderer are never invoked.
Overlap is measured with the EXISTING audit tools (axis-aligned intersection) —
``scripts/audit_pptx_textbox_overlap.py`` (``ov`` / ``boxes``) and
``scripts/audit_pptx_baked_text.py`` (``baked_text_score``).

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_overlay_collision_bug_condition.py -p no:cacheprovider -q

_Bug_Condition: isBugCondition(S) = defectA OR defectB OR defectC (design Bug Condition)_
_Requirements: 1.1, 1.2, 1.3, 1.4, 1.5_
"""
from __future__ import annotations

import os
import sys
import tempfile

# Make ai_engine (repo root) and the audit tools (scripts/) importable.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# Native assembly path (real code under test) — no network.
from ai_engine.native_diagram_pptx import (  # noqa: E402
    build_native_cover,
    build_native_diagram,
)
# Pure decision functions (real code under test) — no network.
from ai_engine.server import _classify_slide_role, _select_render_plan  # noqa: E402
# Pure layout geometry (fixed seam under test) — no network.
import ai_engine.layout_geometry as lg  # noqa: E402

# EXISTING audit measurement functions (reused — same axis-aligned intersection).
import audit_pptx_textbox_overlap as tov  # noqa: E402
import audit_pptx_baked_text as abt  # noqa: E402


# Threshold of "meaningful overlap" — design Bug Condition (10% of smaller box).
THRESHOLD = 0.10
# 16:9 slide so the documented inch coordinates hold (design Examples).
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


def _blank_169_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs):
    # layout[6] = fully blank layout (no placeholders to interfere).
    return prs.slides.add_slide(prs.slide_layouts[6])


def _area(box):
    # box = (head, l, t, w, h)  — same shape audit tools use.
    return max(0.0, box[3]) * max(0.0, box[4])


def _is_badge(box):
    """A number badge: very short purely-numeric label + small near-square box."""
    head = (box[0] or "").strip()
    return head.isdigit() and 1 <= len(head) <= 2


def _worst_pair_overlap(boxes):
    """Return (ratio, info) for the pair with the largest overlap / min(area)."""
    worst = 0.0
    info = None
    for i in range(len(boxes)):
        for j in range(i + 1, len(boxes)):
            ov = tov.ov(boxes[i], boxes[j])
            if ov <= 0.0:
                continue
            amin = min(_area(boxes[i]), _area(boxes[j]))
            if amin <= 0.0:
                continue
            ratio = ov / amin
            if ratio > worst:
                worst = ratio
                info = (boxes[i], boxes[j], ov, amin)
    return worst, info


def _make_baked_text_png(path, w=1280, h=720):
    """Render a full-bleed-ish PNG with text baked into the raster (simulates a
    4K AI background whose Korean labels are rasterised, not editable)."""
    im = Image.new("RGB", (w, h), (28, 42, 74))
    d = ImageDraw.Draw(im)
    line = ("프로젝트 디렉토리 구조 노드 흐름 ROOT SRC DIST BUILD 0123456789 "
            "ABCDEFGH 데이터 정제 적재 분석 파이프라인")
    # Many high-contrast text rows spread across the full width → triggers the
    # baked-text detector (texty rows / runs).
    for k in range(14):
        y = 40 + k * 46
        d.text((30, y), line, fill=(235, 240, 250))
    im.save(path)
    return path


# ──────────────────────────────────────────────────────────────────────────
# (A) defectA — 표지 제목↔부제 수직 겹침
# ──────────────────────────────────────────────────────────────────────────
def test_defect_a_cover_title_subtitle_overlap():
    """build_native_cover 로 긴 제목 표지를 조립 → 제목 박스 하단 > 부제 top →
    텍스트 박스끼리 세로로 겹침(≈10in²). EXPECTED(post-fix): 모든 텍스트 박스
    쌍 겹침 < 10%. 미수정 코드에서는 FAIL 해야 한다 (Req 1.1, 1.2)."""
    prs = _blank_169_presentation()
    slide = _blank_slide(prs)

    drew = build_native_cover(
        slide, prs,
        title="데이터 흐름 구조 정리",            # tlen<=16 → title_pt=46 (>=40) → sub_y=title_y+1.05
        subtitle="조직 전반의 디렉토리 깊이와 데이터 흐름을 한눈에 시각화합니다",
        eyebrow="", date_str="", kpis=None,
    )
    assert drew, "build_native_cover 가 표지를 그리지 못함(전제 실패)"

    boxes = tov.boxes(slide)
    # 표지 텍스트 박스 = 제목 + 부제 두 개여야 한다(전제 확인).
    assert len(boxes) >= 2, f"표지 텍스트 박스가 2개 미만: {boxes}"

    worst, info = _worst_pair_overlap(boxes)
    detail = ""
    if info is not None:
        a, b, ov, amin = info
        detail = (f"  '{a[0]}'@({a[1]},{a[2]} {a[3]}x{a[4]}) ↔ "
                  f"'{b[0]}'@({b[1]},{b[2]} {b[3]}x{b[4]})  "
                  f"overlap={ov:.2f}in² ({worst*100:.0f}% of smaller {amin:.2f}in²)")
    # EXPECTED BEHAVIOUR (fails on unfixed code → proves defectA).
    assert worst < THRESHOLD, (
        "결함 A 재현 — 표지 제목 박스와 부제 박스가 세로로 겹친다 "
        f"(겹침 {worst*100:.0f}% ≥ 임계 {THRESHOLD*100:.0f}%).\n" + detail
    )


# ──────────────────────────────────────────────────────────────────────────
# (C) defectC — 번호 배지↔라벨 박스 겹침
# ──────────────────────────────────────────────────────────────────────────
def test_defect_c_twocol_badge_label_overlap():
    """twocol 6항목 네이티브 카드 조립 → 배지가 라벨 카드 내부에 완전 포함 →
    배지∩라벨 == 배지 면적의 100%. EXPECTED(post-fix): 배지가 라벨 밖 거터,
    겹침 < 배지 면적의 10%. 미수정 코드에서는 FAIL 해야 한다 (Req 1.5)."""
    prs = _blank_169_presentation()
    slide = _blank_slide(prs)

    content = (
        "언어: Python 3.11\n"
        "프레임워크: FastAPI\n"
        "런타임: Electron\n"
        "빌드: webpack\n"
        "테스트: pytest\n"
        "CI: GitHub Actions"
    )
    drew = build_native_diagram(
        slide, "twocol", content,
        region=(0.6, 1.7, 12.1, 5.2), palette=None, title="기술 스택",
    )
    assert drew, "build_native_diagram(twocol) 가 그리지 못함(전제 실패)"

    boxes = tov.boxes(slide)
    badges = [b for b in boxes if _is_badge(b)]
    assert badges, f"번호 배지를 찾지 못함(전제 실패): {[b[0] for b in boxes]}"

    # 각 배지에 대해, 가장 크게 겹치는 (배지가 아닌) 라벨 박스와의 비율(배지 면적 기준).
    worst = 0.0
    info = None
    for bd in badges:
        ba = _area(bd)
        if ba <= 0.0:
            continue
        for lb in boxes:
            if lb is bd or _is_badge(lb):
                continue
            ov = tov.ov(bd, lb)
            if ov <= 0.0:
                continue
            ratio = ov / ba          # design defectC: overlapArea(badge,label) / area(badge)
            if ratio > worst:
                worst = ratio
                info = (bd, lb, ov, ba)

    detail = ""
    if info is not None:
        bd, lb, ov, ba = info
        detail = (f"  배지 '{bd[0]}'@({bd[1]},{bd[2]} {bd[3]}x{bd[4]}) ⊂ "
                  f"라벨 '{lb[0]}'@({lb[1]},{lb[2]} {lb[3]}x{lb[4]})  "
                  f"overlap={ov:.2f}in² ({worst*100:.0f}% of badge {ba:.2f}in²)")
    # EXPECTED BEHAVIOUR (fails on unfixed code → proves defectC).
    assert worst < THRESHOLD, (
        "결함 C 재현 — 번호 배지가 라벨 박스 안에 포개진다 "
        f"(겹침 {worst*100:.0f}% of badge ≥ 임계 {THRESHOLD*100:.0f}%).\n" + detail
    )


# ──────────────────────────────────────────────────────────────────────────
# (B-1) defectB — 구조형이 구워진-텍스트 풀블리드 이미지로 렌더됨
# ──────────────────────────────────────────────────────────────────────────
def test_defect_b1_structural_baked_text_carrier():
    """수정된 결정 seam(``_select_render_plan`` + ``_classify_slide_role``)을 실제로 구동해,
    텍스트가 구워진 풀블리드 배경이 깔린 구조형 슬라이드에서 그 이미지를 본문 캐리어로 쓰지
    않음을 단언한다. EXPECTED(post-fix): ``primary == "NATIVE_SHAPES"`` ∧
    ``vertex_slot == "backdrop"``(손실-0 보존) ∧ ``body_separated is True`` 이고, 시각형
    슬라이드는 구워진-텍스트 시 ``visual`` → ``content`` 로 강등된다(본문 캐리어 차단).
    수정 전 코드(``bg_has_baked_text`` 파라미터/``body_separated`` 키 부재)에서는
    TypeError/KeyError 또는 다른 값으로 FAIL 한다 — 즉 fix 를 진짜로 검증한다 (Req 1.3)."""
    # 0) 전제: 풀블리드 배경에 텍스트가 실제로 구워졌음을 기존 audit 도구로 측정·확인.
    with tempfile.TemporaryDirectory() as td:
        bg_png = _make_baked_text_png(os.path.join(td, "baked_bg.png"))
        pct, lines = abt.baked_text_score(Image.open(bg_png))
    bg_has_baked_text = (pct >= 6.0 or lines >= 6)
    assert bg_has_baked_text, (
        f"전제 실패 — 배경에 텍스트가 구워진 것으로 측정되지 않음(텍스트추정행 {pct}%, 줄 {lines}개)."
    )

    # 1) 수정된 손실-0 결정 함수를 구동: 구조형 + 생성 Vertex 이미지 + 구워진-텍스트 배경.
    plan = _select_render_plan(
        role="structural",
        has_vertex_image=True,
        has_native_diagram=False,
        has_image_file=False,
        has_slide_bg=False,
        html_enabled=False,
        bg_has_baked_text=True,
    )
    # EXPECTED(post-fix) — 구워진-텍스트 이미지를 본문 캐리어로 쓰지 않는다.
    assert plan["primary"] == "NATIVE_SHAPES", (
        "결함 B-1 — 구조형이 편집 네이티브 도형 대신 "
        f"{plan['primary']} 로 렌더된다(구워진-텍스트 이미지가 본문 캐리어)."
    )
    # 손실-0 보존 — 생성 이미지는 폐기되지 않고 backdrop(장식) 슬롯으로 보존.
    assert plan["vertex_slot"] == "backdrop", (
        "결함 B-1 — 생성 Vertex 이미지가 backdrop(손실-0 보존)이 아니라 "
        f"{plan['vertex_slot']!r} 슬롯에 놓인다."
    )
    # 본문/배경 분리 의도가 명시 신호로 전달되어야 한다.
    assert plan["body_separated"] is True, (
        "결함 B-1 — 본문/배경 분리 의도(body_separated)가 신호되지 않는다."
    )

    # 2) 분류기도 함께 구동 — 구워진-텍스트면 시각형 슬라이드를 visual→content 로 강등해
    #    그 이미지를 본문 비주얼 캐리어로 쓰지 못하게 차단한다.
    visual_slide = {
        "title": "팀 협업 워크숍 현장",
        "bullets": ["몰입형 협업 분위기를 담은 사진"],
        "imagePrompt": "사무실에서 노트북으로 함께 일하는 팀의 따뜻하고 자연스러운 사진",
    }
    assert _classify_slide_role(visual_slide, is_cover=False) == "visual", (
        "전제 실패 — 시각형 슬라이드가 visual 로 분류되지 않음(강등 검증 불가)."
    )
    # EXPECTED(post-fix): bg_has_baked_text 시 visual → content 강등.
    assert _classify_slide_role(visual_slide, is_cover=False, bg_has_baked_text=True) == "content", (
        "결함 B-1 — 구워진-텍스트 풀블리드에서도 시각형이 visual 로 남아 본문 캐리어가 된다."
    )


# ──────────────────────────────────────────────────────────────────────────
# (B-2) defectB — 풀블리드 배경 위 본문↔배경 겹침
# ──────────────────────────────────────────────────────────────────────────
def test_defect_b2_body_over_fullbleed_background_overlap():
    """수정된 순수 기하 ``layout_geometry.body_safe_area`` 를 실제로 구동해, 구워진-텍스트
    배경 위 본문을 배경이 점유하지 않는 여백 띠로 분리함을 단언한다. EXPECTED(post-fix):
    분리된 안전 영역 region' 에 대해 ``overlap_area(region', bg) < 0.10 * area(region')``.
    ``body_safe_area`` 가 없던 수정 전에는 import 자체가 불가했으므로 이 호출은 fix 를
    검증한다 (Req 1.4).

    설계 의도(design Fix Implementation §3)대로 배경이 슬라이드를 **완전히 덮지 않는**
    현실적 입력(상단 대부분을 덮되 하단 여백 띠를 남기는 풀블리드-급 배경)으로 구성해,
    본문이 배경 미점유 띠로 이동하는 분리 경로가 실제로 동작하는지 검증한다(배경이 슬라이드를
    완전히 덮어 안전 띠가 없는 극단 케이스는 호출부가 네이티브로 전환하는 별도 설계 경로)."""
    slide = (0.0, 0.0, SLIDE_W_IN, SLIDE_H_IN)
    # 구워진-텍스트 배경이 상단 대부분을 덮되 하단 여백 띠(5.5~7.5)는 남긴다.
    bg = (0.0, 0.0, SLIDE_W_IN, 5.5)
    # 현재(미수정) 코드가 배경 위에 그대로 올리던 본문 region(design Examples).
    desired = (0.5, 1.75, 9.0, 4.95)

    # 전제(결함 재현 조건): 분리 전 본문은 배경과 임계 이상으로 겹친다.
    assert lg.overlap_area(desired, bg) >= THRESHOLD * lg.area(desired), (
        "전제 실패 — desired 본문이 배경과 임계 이상 겹치지 않음(분리할 결함이 없음)."
    )

    # 수정된 fix seam 구동 — 본문을 배경과 분리된 안전 영역으로 산출.
    region = lg.body_safe_area(slide=slide, bg=bg, has_baked_text=True, desired=desired)
    region_area = lg.area(region)
    ov = lg.overlap_area(region, bg)
    ratio = ov / region_area if region_area > 0 else 0.0

    # EXPECTED BEHAVIOUR (fails on unfixed/absent seam → proves defectB-2 is resolved by the fix).
    assert ratio < THRESHOLD, (
        "결함 B-2 — 본문이 풀블리드 배경과 분리되지 않는다 "
        f"(겹침 {ov:.2f}in² = region 의 {ratio*100:.0f}% ≥ 임계 {THRESHOLD*100:.0f}%; "
        f"region={region})."
    )


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
