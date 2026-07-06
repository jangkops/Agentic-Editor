"""엔드투엔드 하위 호환 자동 검증 테스트 (Validates: Requirements 5.2, 9.1, 9.5, 9.6).

이 테스트는 "앱을 직접 구동하지 않고" 백엔드 함수를 직접 호출해
`_force_generate_from_text(...)` → `_resolve_active_template(...)` →
`_tool_generate_pptx(...)` 의 전체 PPTX 생성 경로를 E2E 관점에서 검증한다.

검증 대상(ai_engine/server.py, 구현 절대 수정 금지):
  - _force_generate_from_text(primary_tool, target_files, title, description,
        final_text, project_path, aws_profile, bedrock_user, template_id="")  [async]
        → list of (relative_path, info_dict)
  - _tool_generate_pptx(tool_input, project_path, aws_profile='', bedrock_user='') [async]
  - _resolve_active_template(template_id, store_root) → (path, profile, used)

검증 시나리오:

  1. 하위 호환 (요구사항 5.2):
     templateId를 전혀 전달하지 않은 입력으로 `_force_generate_from_text`를 호출하면
     무템플릿 경로가 선택되어야 한다. 생성 PPTX는
       - 슬라이드 수 = 표지 1 + 섹션 수
       - 레이아웃 매핑 = LAYOUT_MAP {title:0, content:1} (표지 0, 본문 1)
       - 슬라이드 크기 = 16:9 (Inches(13.333) × Inches(7.5), 6858000 EMU)
       - 배경 Tier 단계 = 텍스트 전용(임베드 이미지 0개)
     를 만족하고, `_tool_generate_pptx`에 전달된 입력(inp)에 templatePath/templateId/
     styleProfile 키가 없으며, 그 응답에도 templateId 키가 없어야 한다(기존 형태 보존).

  2. 결정론 (요구사항 5.2):
     같은 입력으로 무템플릿 경로를 2회 호출하면 슬라이드 수·레이아웃 인덱스 시퀀스·
     슬라이드 크기·임베드 이미지 수가 동일해야 한다.

  3. 템플릿 처리 단계 실패 주입에도 유효 PPTX (요구사항 9.1, 9.5, 9.6):
     유효 템플릿을 실제로 등록한 뒤 그 base.pptx를 손상시켜(템플릿 열기 단계 실패 주입)
     동일 입력으로 `_force_generate_from_text`를 호출한다. 그러면
       - `_resolve_active_template`는 base.pptx 파일이 존재하므로 used=True로 해석하여
         inp에 templatePath/templateId/styleProfile을 주입한다(= 템플릿 단계 진입).
       - `_tool_generate_pptx`가 손상 base.pptx 열기에 실패해 무템플릿으로 폴백한다.
     결과적으로 모든 섹션을 포함한 유효 PPTX가 산출되어야 하며(9.5, 9.6), 응답에
     templateId가 없어야 한다(폴백이 일어났음을 증명, 9.1). 슬라이드 수·레이아웃·
     이미지 Tier 단계는 무템플릿 baseline과 동일해야 한다(격리 → 5.2 보존).

네트워크/LLM 차단(앱 직접 구동 없이 함수 호출 기반):
  - 입력 텍스트에 시각 키워드를 넣지 않아 visual_intent=False → Mermaid/matplotlib Tier
    진입 자체가 없다.
  - AE_DISABLE_HTML_SLIDES=1 / AE_DISABLE_MERMAID=1 로 HTML·Mermaid Tier를 끈다.
  - AE_ENABLE_VERTEX_IMAGE 미설정 + vertex_image_module.get_vertex_image_client를
    disabled 스텁으로 patch → AWS Secrets Manager/boto3 접근 차단.
  - server._call_bridge를 None 반환 스텁으로 patch → Electron 브리지 접근 차단.
  - 슬라이드에 imagePrompt/imageFile/slideBackground가 생기지 않으므로 Bedrock
    이미지 호출도 발생하지 않는다.

python-pptx 등 필수 의존성 부재 시 skip. 임시 파일은 TemporaryDirectory로 정리.

실행:
  ai_engine/.venv/bin/python scripts/test_template_end_to_end_backward_compat.py
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile
from unittest.mock import patch

# 레포 루트에서 ai_engine 패키지를 import 가능하게 한다(기존 scripts/ 컨벤션).
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

try:
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception as e:  # pragma: no cover - 의존성 부재 시 skip
    print(f"SKIP: python-pptx import 실패 ({e})")
    sys.exit(0)

try:
    import ai_engine.server as server  # 모듈 객체 — patch 대상 attribute 접근용
    from ai_engine.server import _force_generate_from_text, LAYOUT_MAP
    from ai_engine import template_manager
except Exception as e:  # pragma: no cover - server import 실패 시 skip
    print(f"SKIP: ai_engine import 실패 ({e})")
    sys.exit(0)


# ===========================================================================
# 입력 fixture — 시각 키워드 없는 마크다운(헤더 3개 → 섹션 3개)
# ===========================================================================

# 헤더 3개 → _split_into_sections 가 3개 섹션 반환 → 표지 포함 4 슬라이드.
# 시각 키워드(이미지/그림/다이어그램/차트/visual/diagram ...) 미포함 → visual_intent=False.
_FINAL_TEXT = (
    "# 첫 번째 장\n"
    "\n"
    "- 항목 하나\n"
    "- 항목 둘\n"
    "\n"
    "# 두 번째 장\n"
    "\n"
    "- 항목 셋\n"
    "- 항목 넷\n"
    "\n"
    "# 세 번째 장\n"
    "\n"
    "- 항목 다섯\n"
)

_TITLE = "백워드 호환 검증 자료"
_DESCRIPTION = "분기별 운영 요약 문서"  # 시각 키워드 없음

_SECTION_HEADINGS = ["첫 번째 장", "두 번째 장", "세 번째 장"]
_SECTION_BULLETS = {
    "첫 번째 장": ["항목 하나", "항목 둘"],
    "두 번째 장": ["항목 셋", "항목 넷"],
    "세 번째 장": ["항목 다섯"],
}
_EXPECTED_SLIDE_COUNT = len(_SECTION_HEADINGS) + 1  # +1 표지

# 무템플릿 baseline 레이아웃 인덱스 시퀀스:
#   표지            → _resolve_layout("title", used=False) → LAYOUT_MAP["title"]=0
#   본문 슬라이드 ×3 → sd에 layout 키 없음 → "content" 기본 → LAYOUT_MAP["content"]=1
_EXPECTED_LAYOUT_INDICES = [LAYOUT_MAP["title"]] + [LAYOUT_MAP["content"]] * len(_SECTION_HEADINGS)


# ===========================================================================
# 네트워크 차단 스텁
# ===========================================================================

class _DisabledVertexClient:
    """get_vertex_image_client 대체 — 항상 비활성(Secrets Manager/boto3 미접근)."""

    enabled = False

    def __init__(self, *args, **kwargs):
        self.enabled = False


def _disabled_vertex(*args, **kwargs):
    return _DisabledVertexClient()


def _bridge_unreachable(*args, **kwargs):
    """_call_bridge 대체 — Electron 브리지 미가용(None 반환)."""
    return None


# ===========================================================================
# 실행 헬퍼
# ===========================================================================

def _set_network_block_env():
    """HTML/Mermaid Tier를 끄고 Vertex 옵트인을 해제한다(네트워크 차단 보강)."""
    os.environ["AE_DISABLE_HTML_SLIDES"] = "1"
    os.environ["AE_DISABLE_MERMAID"] = "1"
    os.environ.pop("AE_ENABLE_VERTEX_IMAGE", None)
    os.environ.pop("AE_BEDROCK_HERO_IMAGE", None)
    os.environ.pop("AE_ENABLE_BEDROCK_SLIDE_IMAGES", None)


def _run_force_generate(project_path, template_id="", store_root=None):
    """`_force_generate_from_text`를 네트워크 차단 + 스파이 patch로 실행한다.

    Returns:
        (out, captured)
        - out: _force_generate_from_text 반환 리스트 [(rel, info), ...]
        - captured: [{"inp": dict, "raw": str}, ...] — 실제 _tool_generate_pptx에
          전달된 입력과 그 JSON 응답을 가로챈 기록(E2E 관찰용).
    """
    captured = []
    real_pptx = server._tool_generate_pptx

    async def _spy_pptx(inp, project_path, aws_profile="", bedrock_user=""):
        # inp는 호출부에서 만든 dict — 가로채되 실제 구현을 그대로 호출한다(수정 없음).
        raw = await real_pptx(
            inp, project_path, aws_profile=aws_profile, bedrock_user=bedrock_user
        )
        captured.append({"inp": dict(inp), "raw": raw})
        return raw

    if store_root is not None:
        os.environ["AE_GENERATED_ROOT"] = store_root

    with patch.object(server, "_tool_generate_pptx", new=_spy_pptx), \
         patch.object(server, "_call_bridge", new=_bridge_unreachable), \
         patch("ai_engine.vertex_image_module.get_vertex_image_client", new=_disabled_vertex):
        out = asyncio.run(_force_generate_from_text(
            primary_tool="generate_pptx",
            target_files=[],
            title=_TITLE,
            description=_DESCRIPTION,
            final_text=_FINAL_TEXT,
            project_path=project_path,
            aws_profile="",
            bedrock_user="",
            template_id=template_id,
        ))
    return out, captured


def _locate_pptx(out, captured, project_path):
    """out/captured에서 생성된 .pptx 절대경로를 찾는다."""
    # 1) out 튜플의 absPath 우선
    for rel, info in out:
        if rel.endswith(".pptx"):
            abs_path = info.get("absPath") or ""
            if abs_path and os.path.isfile(abs_path):
                return abs_path
            cand = os.path.join(project_path, rel)
            if os.path.isfile(cand):
                return cand
    # 2) captured 응답의 path로 폴백
    for c in captured:
        try:
            parsed = json.loads(c["raw"])
        except (json.JSONDecodeError, TypeError):
            continue
        rel = parsed.get("path", "")
        if rel.endswith(".pptx"):
            cand = os.path.join(project_path, rel)
            if os.path.isfile(cand):
                return cand
    return None


def _layout_indices(prs):
    """각 슬라이드의 slide_layout이 prs.slide_layouts에서 차지하는 인덱스 리스트."""
    layouts = list(prs.slide_layouts)
    return [layouts.index(s.slide_layout) for s in prs.slides]


def _picture_count(prs):
    """프레젠테이션 전체에서 임베드된 PICTURE 셰이프 개수(배경 Tier 산출물 지표)."""
    n = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    n += 1
            except Exception:
                # shape_type 접근 실패는 그림 아님으로 간주
                continue
    return n


def _collect_text(prs):
    """모든 슬라이드 텍스트 프레임 내용을 하나의 문자열로 모은다."""
    chunks = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                chunks.append(shape.text_frame.text)
    return "\n".join(chunks)


def _make_valid_base_pptx(path):
    """등록용 유효 base.pptx를 생성한다(python-pptx 기본 템플릿)."""
    prs = Presentation()
    # 표지 한 장 추가 — 유효한 .pptx면 충분(register는 레이아웃/테마만 사용).
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(path)


# ===========================================================================
# 1. 하위 호환 — 무템플릿 E2E (요구사항 5.2)
# ===========================================================================

def test_e2e_no_template_backward_compat():
    """templateId 미전달 → 무템플릿 경로, baseline과 슬라이드 수·레이아웃·크기·Tier 동일."""
    _set_network_block_env()
    with tempfile.TemporaryDirectory() as tmp:
        work = os.path.join(tmp, "work")
        store = os.path.join(tmp, "store")
        os.makedirs(work, exist_ok=True)
        os.makedirs(store, exist_ok=True)

        out, captured = _run_force_generate(work, template_id="", store_root=store)

        # _tool_generate_pptx 가 정확히 1회 호출됐고, 그 입력에 템플릿 키가 없어야 한다(5.2).
        assert len(captured) == 1, f"_tool_generate_pptx 호출 횟수 비정상: {len(captured)}"
        inp = captured[0]["inp"]
        for k in ("templatePath", "templateId", "styleProfile"):
            assert k not in inp, f"무템플릿인데 inp에 '{k}'가 주입됨: {list(inp.keys())}"

        # 내부 응답에 templateId 키가 없어야 한다(기존 응답 형태 보존, 5.2).
        parsed = json.loads(captured[0]["raw"])
        assert "error" not in parsed, f"무템플릿 생성 에러: {parsed}"
        assert "templateId" not in parsed, f"무템플릿 응답에 templateId 포함: {parsed}"
        assert parsed.get("slideCount") == _EXPECTED_SLIDE_COUNT, (
            f"응답 slideCount 불일치: got {parsed.get('slideCount')}, "
            f"expected {_EXPECTED_SLIDE_COUNT}"
        )

        # 생성 .pptx를 찾아 재파싱 검증.
        abs_path = _locate_pptx(out, captured, work)
        assert abs_path is not None, f"생성된 .pptx를 찾지 못함. out={out}"
        prs = Presentation(abs_path)

        # 슬라이드 수 = 표지 1 + 섹션 수
        assert len(prs.slides) == _EXPECTED_SLIDE_COUNT, (
            f"파싱 슬라이드 수 {len(prs.slides)} != 기대 {_EXPECTED_SLIDE_COUNT}"
        )

        # 레이아웃 매핑 = LAYOUT_MAP (표지 0, 본문 1)
        indices = _layout_indices(prs)
        assert indices == _EXPECTED_LAYOUT_INDICES, (
            f"레이아웃 인덱스 시퀀스 불일치: got {indices}, "
            f"expected {_EXPECTED_LAYOUT_INDICES}"
        )

        # 슬라이드 크기 16:9 (13.333 × 7.5 in, 7.5in == 6858000 EMU)
        assert prs.slide_width == Inches(13.333), (
            f"slide_width EMU 불일치: {prs.slide_width} != {int(Inches(13.333))}"
        )
        assert prs.slide_height == Inches(7.5), (
            f"slide_height EMU 불일치: {prs.slide_height} != {int(Inches(7.5))}"
        )
        assert int(prs.slide_height) == 6858000, int(prs.slide_height)

        # 배경 Tier 단계 = 텍스트 전용(임베드 이미지 0개) — visual_intent=False 경로 보존
        assert _picture_count(prs) == 0, (
            f"무템플릿 텍스트 전용인데 임베드 이미지 발견: {_picture_count(prs)}개"
        )

        # 모든 섹션 제목/본문이 텍스트로 포함됐는지 확인(편집 가능 텍스트, 9.5 정신)
        joined = _collect_text(prs)
        assert _TITLE in joined, f"표지 제목 누락: {_TITLE}"
        for heading in _SECTION_HEADINGS:
            assert heading in joined, f"섹션 제목 누락: {heading}"
            for bullet in _SECTION_BULLETS[heading]:
                assert bullet in joined, f"본문 누락: {bullet}"

    print("  1. 무템플릿 E2E → 슬라이드수·레이아웃·16:9·텍스트전용 baseline 동일  OK")


# ===========================================================================
# 2. 결정론 — 같은 입력 2회 호출 동일성 (요구사항 5.2)
# ===========================================================================

def test_e2e_no_template_deterministic():
    """같은 입력으로 무템플릿 경로 2회 호출 → 슬라이드수·레이아웃·크기·이미지수 동일."""
    _set_network_block_env()
    results = []
    for _ in range(2):
        with tempfile.TemporaryDirectory() as tmp:
            work = os.path.join(tmp, "work")
            store = os.path.join(tmp, "store")
            os.makedirs(work, exist_ok=True)
            os.makedirs(store, exist_ok=True)

            out, captured = _run_force_generate(work, template_id="", store_root=store)
            abs_path = _locate_pptx(out, captured, work)
            assert abs_path is not None, f"생성된 .pptx를 찾지 못함. out={out}"
            prs = Presentation(abs_path)
            results.append({
                "slides": len(prs.slides),
                "layouts": _layout_indices(prs),
                "width": int(prs.slide_width),
                "height": int(prs.slide_height),
                "pictures": _picture_count(prs),
            })

    assert results[0] == results[1], (
        f"무템플릿 경로가 결정론적이지 않음:\n  run1={results[0]}\n  run2={results[1]}"
    )
    # 기대값과도 일치(회귀 방지)
    assert results[0]["slides"] == _EXPECTED_SLIDE_COUNT, results[0]
    assert results[0]["layouts"] == _EXPECTED_LAYOUT_INDICES, results[0]
    assert results[0]["pictures"] == 0, results[0]
    print("  2. 무템플릿 E2E 2회 호출 결정론적 동일(슬라이드/레이아웃/크기/이미지)     OK")


# ===========================================================================
# 3. 템플릿 단계 실패 주입에도 유효 PPTX (요구사항 9.1, 9.5, 9.6)
# ===========================================================================

def test_e2e_template_stage_failure_injection_still_valid():
    """유효 템플릿 등록 후 base.pptx 손상 주입 → E2E 폴백, 모든 슬라이드 유효 PPTX."""
    _set_network_block_env()
    with tempfile.TemporaryDirectory() as tmp:
        work = os.path.join(tmp, "work")
        store = os.path.join(tmp, "store")
        os.makedirs(work, exist_ok=True)
        os.makedirs(store, exist_ok=True)
        os.environ["AE_GENERATED_ROOT"] = store

        # (a) 유효 base.pptx 생성 → register_template로 실제 등록.
        src_pptx = os.path.join(tmp, "src_template.pptx")
        _make_valid_base_pptx(src_pptx)
        reg = template_manager.register_template(src_pptx, "E2E 손상 주입 템플릿", store_root=store)
        assert isinstance(reg, dict) and "templateId" in reg, f"템플릿 등록 실패: {reg}"
        template_id = reg["templateId"]

        # (b) 템플릿 처리 단계 실패 주입 — 저장된 base.pptx를 손상 바이트로 덮어쓴다.
        #     metadata.json/style_profile.json은 그대로 두어 _resolve_active_template가
        #     used=True로 해석(템플릿 단계 진입)하도록 한다. 실제 열기에서 실패가 발생한다.
        base_pptx = os.path.join(store, "templates", template_id, "base.pptx")
        assert os.path.isfile(base_pptx), f"등록된 base.pptx 부재: {base_pptx}"
        with open(base_pptx, "wb") as f:
            f.write(b"corrupted-not-a-zip-pptx \x00\x01\x02\x03")

        # (c) 손상 템플릿을 활성 템플릿으로 지정해 E2E 호출.
        out, captured = _run_force_generate(work, template_id=template_id, store_root=store)

        assert len(captured) == 1, f"_tool_generate_pptx 호출 횟수 비정상: {len(captured)}"
        inp = captured[0]["inp"]

        # 템플릿 단계에 실제로 진입했음을 증명 — inp에 templateId/templatePath가 주입됨.
        assert inp.get("templateId") == template_id, (
            f"템플릿 단계 미진입(templateId 미주입): {list(inp.keys())}"
        )
        assert inp.get("templatePath"), f"templatePath 미주입: {list(inp.keys())}"
        assert inp["templatePath"].endswith("base.pptx"), inp["templatePath"]

        # 폴백이 일어나 응답에 templateId가 없어야 한다(템플릿 열기 실패 격리, 9.1).
        parsed = json.loads(captured[0]["raw"])
        assert "error" not in parsed, f"폴백 실패 — 에러 반환: {parsed}"
        assert "templateId" not in parsed, (
            f"손상 템플릿 열기 실패인데 응답에 templateId 포함(폴백 미발생): {parsed}"
        )
        assert parsed.get("slideCount") == _EXPECTED_SLIDE_COUNT, (
            f"응답 slideCount 불일치: got {parsed.get('slideCount')}, "
            f"expected {_EXPECTED_SLIDE_COUNT}"
        )

        # 생성 PPTX 재파싱 — 모든 슬라이드 포함 유효 산출 확인(9.5, 9.6).
        abs_path = _locate_pptx(out, captured, work)
        assert abs_path is not None, f"생성된 .pptx를 찾지 못함. out={out}"
        prs = Presentation(abs_path)

        assert len(prs.slides) == _EXPECTED_SLIDE_COUNT, (
            f"파싱 슬라이드 수 {len(prs.slides)} != 기대 {_EXPECTED_SLIDE_COUNT}"
        )
        # 격리 → 무템플릿 baseline과 레이아웃 매핑·이미지 Tier 동일(5.2 보존).
        assert _layout_indices(prs) == _EXPECTED_LAYOUT_INDICES, (
            f"폴백 후 레이아웃 인덱스 불일치: {_layout_indices(prs)} != {_EXPECTED_LAYOUT_INDICES}"
        )
        assert _picture_count(prs) == 0, (
            f"폴백 후 임베드 이미지 발견(Tier 단계 불일치): {_picture_count(prs)}개"
        )
        # 16:9 baseline 크기로 폴백
        assert prs.slide_width == Inches(13.333) and prs.slide_height == Inches(7.5), (
            f"폴백 슬라이드 크기 불일치: {prs.slide_width}×{prs.slide_height}"
        )

        # 모든 섹션 텍스트 포함(9.5, 9.6)
        joined = _collect_text(prs)
        assert _TITLE in joined, f"표지 제목 누락: {_TITLE}"
        for heading in _SECTION_HEADINGS:
            assert heading in joined, f"섹션 제목 누락: {heading}"
            for bullet in _SECTION_BULLETS[heading]:
                assert bullet in joined, f"본문 누락: {bullet}"

    print("  3. 템플릿 단계 실패 주입 → E2E 폴백 + 모든 슬라이드 유효 PPTX            OK")


# ===========================================================================
# 4. 무템플릿 baseline vs 실패 주입 산출 동일성 (요구사항 5.2 + 9.6)
# ===========================================================================

def test_e2e_baseline_equals_failure_injection_structure():
    """무템플릿 baseline과 템플릿 실패 주입 산출의 구조(슬라이드수/레이아웃/Tier)가 동일."""
    _set_network_block_env()

    def _structure_no_template():
        with tempfile.TemporaryDirectory() as tmp:
            work = os.path.join(tmp, "work")
            store = os.path.join(tmp, "store")
            os.makedirs(work, exist_ok=True)
            os.makedirs(store, exist_ok=True)
            out, captured = _run_force_generate(work, template_id="", store_root=store)
            abs_path = _locate_pptx(out, captured, work)
            assert abs_path is not None, "무템플릿 .pptx 부재"
            prs = Presentation(abs_path)
            return {
                "slides": len(prs.slides),
                "layouts": _layout_indices(prs),
                "width": int(prs.slide_width),
                "height": int(prs.slide_height),
                "pictures": _picture_count(prs),
            }

    def _structure_failure_injection():
        with tempfile.TemporaryDirectory() as tmp:
            work = os.path.join(tmp, "work")
            store = os.path.join(tmp, "store")
            os.makedirs(work, exist_ok=True)
            os.makedirs(store, exist_ok=True)
            os.environ["AE_GENERATED_ROOT"] = store
            src_pptx = os.path.join(tmp, "src.pptx")
            _make_valid_base_pptx(src_pptx)
            reg = template_manager.register_template(src_pptx, "구조 비교 템플릿", store_root=store)
            assert "templateId" in reg, f"등록 실패: {reg}"
            tid = reg["templateId"]
            base_pptx = os.path.join(store, "templates", tid, "base.pptx")
            with open(base_pptx, "wb") as f:
                f.write(b"broken \x00\xff")
            out, captured = _run_force_generate(work, template_id=tid, store_root=store)
            abs_path = _locate_pptx(out, captured, work)
            assert abs_path is not None, "실패주입 .pptx 부재"
            prs = Presentation(abs_path)
            return {
                "slides": len(prs.slides),
                "layouts": _layout_indices(prs),
                "width": int(prs.slide_width),
                "height": int(prs.slide_height),
                "pictures": _picture_count(prs),
            }

    base = _structure_no_template()
    inj = _structure_failure_injection()
    assert base == inj, (
        f"무템플릿 baseline과 실패 주입 산출 구조 불일치:\n"
        f"  baseline={base}\n  injection={inj}"
    )
    print("  4. 무템플릿 baseline ≡ 템플릿 실패 주입 산출 구조(5.2 보존)             OK")


def main():
    print("=== 엔드투엔드 하위 호환 자동 검증 (요구사항 5.2, 9.1, 9.5, 9.6) ===")
    test_e2e_no_template_backward_compat()
    test_e2e_no_template_deterministic()
    test_e2e_template_stage_failure_injection_still_valid()
    test_e2e_baseline_equals_failure_injection_structure()
    print("모든 엔드투엔드 하위 호환 케이스 통과.")


if __name__ == "__main__":
    main()
