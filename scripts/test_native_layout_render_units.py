"""단위/엣지 테스트 — native_layout_renderer / parity_scorer / audit_pptx_native_density.

작업 17.2 (spec: pptx-native-density-render). 검증 범위:

  - 각 ``emit_*`` 가 기대 도형 종류(텍스트박스/ROUNDED_RECTANGLE/OVAL/RECTANGLE/
    PICTURE-or-placeholder)와 design_tokens 색을 생성하는지 (Req 1.3, 5.4).
  - ``apply_tokens_to_run`` 이 토큰 색·폰트 크기를 run 에 적용하는지 (Req 5.4).
  - ``parity_scorer.score(html, category)`` 회귀 — 대표 HTML 로 passed/missing 확인
    (Req 5.2, 5.3).
  - 과밀 입력 → ``finalize_placement`` 가 ``OverlapError`` 발생 (Req 2.5).
  - 빈/미지원 카테고리 → ``score`` / ``score_native_slide`` ``ValueError`` (Req 5.5).
  - 정적 제약 — 렌더러/기하 모듈 소스에 네트워크/게이트웨이 직접 import 없음, 이미지
    생성은 vertex_image_module 단일 모듈만 경유 (Req 10.1, 10.2).
  - 기존 네이티브 다이어그램 경로 보존 — build_native_diagram/build_native_cover
    import 가능 (Req 9.3).
  - Vertex 호출 mock — ``_get_vertex_client`` monkeypatch 로 enabled=True mock 시
    ``maybe_generate_decorative`` 가 단일 모듈 경유함 확인 (실제 호출 금지, Req 11.1).

전부 hermetic — 순수 Python, 네트워크 0, python-pptx in-memory.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_native_layout_render_units.py -p no:cacheprovider -q
"""
from __future__ import annotations

import os
import sys
import base64
import inspect

# repo 루트 + ai_engine 를 import 경로에 추가.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "ai_engine"))

import pytest  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.native_layout_renderer as nlr  # noqa: E402
from ai_engine.native_layout_renderer import (  # noqa: E402
    PlacedShape,
    OverlapError,
    emit_title,
    emit_section_header_bar,
    emit_contact_box,
    emit_note_callout,
    emit_numbered_list,
    emit_card_grid,
    emit_figure_slot,
    apply_tokens_to_run,
    finalize_placement,
    maybe_generate_decorative,
)
from ai_engine.slide_templates import design_tokens_for_profile  # noqa: E402

import scripts.parity_scorer as ps  # noqa: E402


# ---------------------------------------------------------------------------
# 공통 픽스처/헬퍼
# ---------------------------------------------------------------------------

# 1x1 PNG (손실-0 임베드 테스트용 최소 유효 이미지).
_PNG_1x1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mNk"
    "+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
)


@pytest.fixture()
def tokens():
    return design_tokens_for_profile(None)


@pytest.fixture()
def slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    return prs.slides.add_slide(blank)


def _hex(rgb) -> str:
    """RGBColor → 'RRGGBB' 대문자."""
    return str(rgb).upper()


def _autoshapes(slide, name_substr: str):
    """slide 에서 auto_shape_type 이름에 name_substr 를 포함하는 오토셰이프 목록."""
    out = []
    for shp in slide.shapes:
        try:
            if shp.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
                continue
            at = shp.auto_shape_type
            nm = (getattr(at, "name", None) or str(at)).upper()
        except Exception:
            continue
        if name_substr in nm:
            out.append(shp)
    return out


def _rounded_rects(slide):
    """ROUNDED_RECTANGLE 만(일반 RECTANGLE 제외)."""
    return [s for s in _autoshapes(slide, "ROUNDED_RECTANGLE")]


def _plain_rects(slide):
    """일반 RECTANGLE 만(ROUNDED 제외)."""
    out = []
    for s in _autoshapes(slide, "RECTANGLE"):
        at = s.auto_shape_type
        nm = (getattr(at, "name", None) or str(at)).upper()
        if "ROUNDED" not in nm:
            out.append(s)
    return out


def _fill_hex(shp):
    return _hex(shp.fill.fore_color.rgb)


# ===========================================================================
# emit_* 도형 종류 + design_tokens 색 (Req 1.3, 5.4)
# ===========================================================================


def test_emit_title_is_single_editable_textbox(slide, tokens):
    before = len(slide.shapes)
    ph = emit_title(slide, "전략 로드맵", tokens, (0.7, 0.5, 11.9, 1.0))
    assert len(slide.shapes) == before + 1
    assert ph.role == "title"
    assert ph.has_text is True
    tb = slide.shapes[-1]
    assert tb.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
    assert "전략 로드맵" in tb.text_frame.text


def test_emit_section_header_bar_shapes_and_token_colors(slide, tokens):
    placed = emit_section_header_bar(slide, 2, "도입 배경", tokens, (0.7, 2.6, 11.9, 1.1))
    roles = {p.role for p in placed}
    assert {"section_bar", "badge", "section_title"} <= roles

    # 다크 막대 = ROUNDED_RECTANGLE, text_dark 토큰 색.
    bars = _rounded_rects(slide)
    assert len(bars) >= 1
    assert any(_fill_hex(b) == _hex_token(tokens, "text_dark") for b in bars)

    # 번호 배지 = OVAL, accent 토큰 색, 번호 텍스트.
    ovals = _autoshapes(slide, "OVAL")
    assert len(ovals) == 1
    assert _fill_hex(ovals[0]) == _hex_token(tokens, "accent")
    assert ovals[0].text_frame.text.strip() == "2"

    # 제목 텍스트는 편집가능 텍스트박스로 존재.
    assert any(
        s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX and "도입 배경" in s.text_frame.text
        for s in slide.shapes
    )


def test_emit_contact_box_shapes_and_accent_bar(slide, tokens):
    placed = emit_contact_box(slide, {"email": "a@b.com", "tel": "010"}, tokens,
                              (0.7, 2.0, 5.0, 2.0))
    roles = [p.role for p in placed]
    assert "contact_box" in roles and "accent_bar" in roles and "contact" in roles

    # 컨테이너 = ROUNDED_RECTANGLE, 좌측 액센트 바 = 일반 RECTANGLE(accent 색).
    assert len(_rounded_rects(slide)) >= 1
    bars = _plain_rects(slide)
    assert len(bars) >= 1
    assert any(_fill_hex(b) == _hex_token(tokens, "accent") for b in bars)
    assert any("a@b.com" in s.text_frame.text for s in slide.shapes
              if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX)


def test_emit_note_callout_has_notice_label_and_accent_border(slide, tokens):
    placed = emit_note_callout(slide, "마감 임박", tokens, (0.7, 2.0, 6.0, 1.5))
    roles = [p.role for p in placed]
    assert "note_box" in roles and "note_border" in roles and "note" in roles

    # 라운드 박스 + 좌측 보더(accent) + NOTICE 라벨 텍스트.
    assert len(_rounded_rects(slide)) >= 1
    bars = _plain_rects(slide)
    assert any(_fill_hex(b) == _hex_token(tokens, "accent") for b in bars)
    all_text = " ".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    assert "NOTICE" in all_text.upper()
    assert "마감 임박" in all_text


def test_emit_numbered_list_badges_and_text_no_overlap(slide, tokens):
    items = ["계획", "실행", "검토"]
    placed = emit_numbered_list(slide, items, tokens, (0.7, 1.9, 11.9, 4.0))
    # 항목마다 배지(OVAL) + 본문(body).
    ovals = _autoshapes(slide, "OVAL")
    assert len(ovals) == 3
    assert all(_fill_hex(o) == _hex_token(tokens, "accent") for o in ovals)
    bodies = [p for p in placed if p.role == "body"]
    assert len(bodies) == 3

    # 거터 배치 — 각 행의 배지 우단 <= 본문 좌단 (겹침 0).
    badges = [p for p in placed if p.role == "badge"]
    for bd, bo in zip(badges, bodies):
        badge_right = bd.rect[0] + bd.rect[2]
        body_left = bo.rect[0]
        assert badge_right <= body_left + 1e-6


def test_emit_card_grid_rounded_cards_with_card_bg_token(slide, tokens):
    cards = [{"title": "A", "description": "alpha"},
             {"title": "B", "description": "beta"},
             {"title": "C", "description": "gamma"}]
    placed = emit_card_grid(slide, cards, tokens, (0.7, 1.9, 11.9, 4.5))
    card_shapes = _rounded_rects(slide)
    assert len(card_shapes) == 3
    assert all(_fill_hex(c) == _hex_token(tokens, "card_bg") for c in card_shapes)
    assert all(p.role == "card" and p.has_text for p in placed)


def test_emit_figure_slot_placeholder_when_no_image(slide):
    ph = emit_figure_slot(slide, (1.0, 1.0, 4.0, 3.0), image_path=None)
    assert ph.role == "figure"
    assert ph.has_text is False  # 콘텐츠 텍스트 베이크 금지.
    rects = _plain_rects(slide)
    assert len(rects) == 1
    # 중립 회색 플레이스홀더(콘텐츠 색 아님).
    assert _fill_hex(rects[0]) == "F0F2F5"


def test_emit_figure_slot_embeds_picture_with_image(slide, tmp_path):
    img = tmp_path / "deco.png"
    img.write_bytes(_PNG_1x1)
    before = len(slide.shapes)
    ph = emit_figure_slot(slide, (1.0, 1.0, 4.0, 3.0), image_path=str(img))
    assert ph.role == "figure" and ph.has_text is False
    assert len(slide.shapes) == before + 1
    assert slide.shapes[-1].shape_type == MSO_SHAPE_TYPE.PICTURE


# ===========================================================================
# apply_tokens_to_run — 토큰 색/폰트 크기 적용 (Req 5.4)
# ===========================================================================


def _hex_token(tokens, key) -> str:
    """tokens[key] ('#RRGGBB') → 'RRGGBB' 대문자."""
    return str(tokens[key]).lstrip("#").upper()


def test_apply_tokens_to_run_title_uses_token_color_and_large_font(slide, tokens):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "제목"
    apply_tokens_to_run(run, tokens, "title")
    assert run.font.size == Pt(30)
    assert run.font.bold is True
    assert _hex(run.font.color.rgb) == _hex_token(tokens, "text_dark")


def test_apply_tokens_to_run_typography_hierarchy_title_gt_body(slide, tokens):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    p = tb.text_frame.paragraphs[0]
    r_title = p.add_run(); r_title.text = "T"
    r_body = tb.text_frame.add_paragraph().add_run(); r_body.text = "b"
    apply_tokens_to_run(r_title, tokens, "title")
    apply_tokens_to_run(r_body, tokens, "body")
    assert r_title.font.size > r_body.font.size  # 타이포 계층(제목 > 본문).


def test_apply_tokens_to_run_badge_uses_accent_or_light_token(slide, tokens):
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
    run = tb.text_frame.paragraphs[0].add_run()
    run.text = "1"
    apply_tokens_to_run(run, tokens, "badge")
    # badge role 은 text_light 토큰 색(흰 텍스트) 사용.
    assert _hex(run.font.color.rgb) == _hex_token(tokens, "text_light")


# ===========================================================================
# parity_scorer.score 회귀 — 대표 HTML passed/missing (Req 5.2, 5.3)
# ===========================================================================


def _cover_html(marker_names) -> str:
    """COVER_CHECKLIST 의 일부 마커를 포함하는 대표 HTML 조립."""
    by_name = dict(ps.COVER_CHECKLIST)
    parts = [f'<div {by_name[n]}>x</div>' for n in marker_names]
    return "<section>" + "".join(parts) + "</section>"


def _body_html(marker_names) -> str:
    by_name = dict(ps.BODY_CHECKLIST)
    parts = [f'<div {by_name[n]}>x</div>' for n in marker_names]
    return "<section>" + "".join(parts) + "</section>"


def test_score_cover_passed_with_reference_markers():
    names = ["icon_badge", "notice_chip", "accent_head",
             "step_grid", "accent_bar", "corner_glow"]  # 6 >= ref 6
    res = ps.score(_cover_html(names), "cover")
    assert res["density_score"] == 6
    assert res["reference_score"] == 6
    assert res["passed"] is True
    assert res["missing"] == ["footer"]  # 유일 미충족 항목.


def test_score_body_passed_and_missing_reported():
    names = ["section_header", "contact_box", "note_callout",
             "link_chip", "numbered_item", "notice_tab"]  # 6 >= ref 6
    res = ps.score(_body_html(names), "body")
    assert res["passed"] is True
    assert set(res["missing"]) == {"slide_footer", "figure_slot"}


def test_score_body_failed_when_below_reference():
    names = ["section_header", "contact_box"]  # 2 < ref 6
    res = ps.score(_body_html(names), "body")
    assert res["density_score"] == 2
    assert res["passed"] is False
    assert len(res["missing"]) == 6


def test_score_native_slide_matches_score_dict_shape(slide, tokens):
    # 네이티브 어댑터도 score() 와 동일 키 집합을 반환해야 한다(형식 회귀).
    emit_title(slide, "본문 제목", tokens, (0.7, 0.5, 11.9, 1.0))
    res = ps.score_native_slide(slide, "body")
    assert set(res.keys()) == {
        "category", "density_score", "reference_score",
        "total", "passed", "items", "missing",
    }
    # present ∪ missing == 전체, 상호배타(일관성).
    present = {it["name"] for it in res["items"] if it["present"]}
    missing = set(res["missing"])
    allnames = {n for n, _ in ps.BODY_CHECKLIST}
    assert present | missing == allnames
    assert present & missing == set()


# ===========================================================================
# finalize_placement — 과밀 입력 OverlapError (Req 2.5)
# ===========================================================================


def test_finalize_placement_raises_overlaperror_on_overcrowded():
    # 슬라이드 전체를 덮는 텍스트 셰이프 3개 → 보정 불가 → OverlapError.
    full = (0.0, 0.0, nlr.SLIDE_W_IN, nlr.SLIDE_H_IN)
    placed = [
        PlacedShape(role="body", rect=full, has_text=True, text=f"blk{i}", z=2)
        for i in range(3)
    ]
    with pytest.raises(OverlapError) as ei:
        finalize_placement(placed, slide_id="S1")
    err = ei.value
    assert err.slide_id == "S1"
    assert len(err.pairs) >= 1  # 위반 셰이프 쌍 보고.


def test_finalize_placement_noop_for_clean_input():
    # 이미 경계 안·겹침 없는 입력 → 좌표 불변(no-op).
    placed = [
        PlacedShape(role="title", rect=(0.7, 0.5, 11.9, 1.0), has_text=True, text="t", z=5),
        PlacedShape(role="body", rect=(0.7, 2.0, 5.0, 3.0), has_text=True, text="b", z=2),
    ]
    out = finalize_placement(placed)
    assert [p.rect for p in out] == [(0.7, 0.5, 11.9, 1.0), (0.7, 2.0, 5.0, 3.0)]


def test_finalize_placement_dedups_identical_titles():
    placed = [
        PlacedShape(role="title", rect=(0.7, 0.5, 5.0, 1.0), has_text=True, text="개요", z=5),
        PlacedShape(role="title", rect=(0.7, 3.0, 5.0, 1.0), has_text=True, text="개요", z=5),
    ]
    out = finalize_placement(placed)
    titles = [p for p in out if p.role in nlr.TITLE_ROLES]
    assert len(titles) == 1  # 동일 제목 1개만 유지(Req 4.2).


# ===========================================================================
# 빈/미지원 카테고리 ValueError (Req 5.5)
# ===========================================================================


@pytest.mark.parametrize("html,category", [
    (None, "cover"),
    ("", "body"),
])
def test_score_value_error_on_empty_html(html, category):
    with pytest.raises(ValueError):
        ps.score(html, category)


@pytest.mark.parametrize("category", ["timeline", "", "Cover", "table"])
def test_score_value_error_on_unsupported_category(category):
    with pytest.raises(ValueError):
        ps.score("<div>x</div>", category)


def test_score_native_slide_value_error_on_none_slide():
    with pytest.raises(ValueError):
        ps.score_native_slide(None, "body")


@pytest.mark.parametrize("category", ["timeline", "", "Body", "diagram"])
def test_score_native_slide_value_error_on_unsupported_category(slide, category):
    with pytest.raises(ValueError):
        ps.score_native_slide(slide, category)


# ===========================================================================
# 정적 제약 — 네트워크/게이트웨이 import 없음, 이미지 생성 단일 모듈 (Req 10.1, 10.2)
# ===========================================================================


def _module_source(mod) -> str:
    return inspect.getsource(mod)


@pytest.mark.parametrize("modname", [
    "ai_engine.native_layout_renderer",
    "ai_engine.layout_geometry",
])
def test_no_network_or_gateway_imports(modname):
    import importlib
    mod = importlib.import_module(modname)
    src = _module_source(mod)
    # 네트워크/게이트웨이 직접 import 금지(렌더러/기하는 순수 렌더+기하).
    forbidden = [
        "import httpx", "import requests",
        "from httpx", "from requests",
        "import gateway_module", "from gateway_module",
        "import boto3",
    ]
    for token in forbidden:
        assert token not in src, f"{modname} 에 금지된 import 발견: {token!r}"


def test_image_generation_routes_through_vertex_image_module_only():
    # 이미지 생성은 _get_vertex_client(=vertex_image_module 위임)만 경유해야 한다.
    src_client = inspect.getsource(nlr._get_vertex_client)
    assert "vertex_image_module" in src_client
    assert "get_vertex_image_client" in src_client

    # maybe_generate_decorative 는 _get_vertex_client 를 통해서만 클라이언트를 얻는다
    # (다른 이미지 API 클라이언트를 직접 생성하지 않음).
    src_gen = inspect.getsource(maybe_generate_decorative)
    assert "_get_vertex_client" in src_gen


# ===========================================================================
# 기존 네이티브 다이어그램 경로 보존 (Req 9.3)
# ===========================================================================


def test_existing_native_diagram_path_importable():
    from ai_engine.native_diagram_pptx import build_native_diagram, build_native_cover
    assert callable(build_native_diagram)
    assert callable(build_native_cover)
    # 렌더러도 동일 빌더를 재사용 배선해 두었는지 확인(회귀 방지).
    assert callable(nlr.build_native_diagram)
    assert callable(nlr.build_native_cover)


# ===========================================================================
# Vertex 호출 mock — 단일 모듈 경유 (Req 11.1) — 실제 호출 금지
# ===========================================================================


class _FakeVertexClient:
    """네트워크 없는 mock Vertex 클라이언트. generate 호출 수를 센다."""

    def __init__(self, enabled=True, image_b64=None, result=None):
        self.enabled = enabled
        self._image_b64 = image_b64
        self._result = result
        self.calls = 0

    def generate(self, prompt, model_class=None, aspect_ratio="16:9", **kw):
        self.calls += 1
        if self._result is not None:
            return self._result
        return {"images": [self._image_b64], "model": "mock"}


def test_maybe_generate_decorative_off_returns_none(monkeypatch):
    monkeypatch.delenv("AE_ENABLE_VERTEX_IMAGE", raising=False)
    fake = _FakeVertexClient(enabled=True)
    monkeypatch.setattr(nlr, "_get_vertex_client", lambda **kw: fake)
    out = maybe_generate_decorative("abstract bg", "background")
    assert out is None
    assert fake.calls == 0  # 옵트인 OFF → 생성 호출조차 안 함.


def test_maybe_generate_decorative_mock_on_routes_single_module(monkeypatch, tmp_path):
    monkeypatch.setenv("AE_ENABLE_VERTEX_IMAGE", "1")
    b64 = base64.b64encode(_PNG_1x1).decode("ascii")
    fake = _FakeVertexClient(enabled=True, image_b64=b64)
    monkeypatch.setattr(nlr, "_get_vertex_client", lambda **kw: fake)

    out = maybe_generate_decorative("abstract bg, no text", "background",
                                    out_dir=str(tmp_path))
    assert out is not None and os.path.exists(out)
    assert fake.calls == 1  # 단일 모듈(_get_vertex_client) 경유 1회.
    # 손실-0: 저장된 바이트 == 원본 디코드 바이트(재인코딩 없음).
    with open(out, "rb") as f:
        assert f.read() == _PNG_1x1


def test_maybe_generate_decorative_disabled_client_returns_none(monkeypatch):
    monkeypatch.setenv("AE_ENABLE_VERTEX_IMAGE", "1")
    fake = _FakeVertexClient(enabled=False)  # 자격증명 부재 → enabled=False.
    monkeypatch.setattr(nlr, "_get_vertex_client", lambda **kw: fake)
    out = maybe_generate_decorative("abstract bg", "background")
    assert out is None
    assert fake.calls == 0


def test_maybe_generate_decorative_error_result_returns_none(monkeypatch):
    monkeypatch.setenv("AE_ENABLE_VERTEX_IMAGE", "1")
    fake = _FakeVertexClient(enabled=True, result={"error": "boom"})
    monkeypatch.setattr(nlr, "_get_vertex_client", lambda **kw: fake)
    out = maybe_generate_decorative("abstract bg", "background")
    assert out is None  # 생성 실패 → None, 예외 전파 금지.
    assert fake.calls == 1
