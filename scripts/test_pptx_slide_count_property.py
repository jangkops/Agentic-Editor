"""Property 8: PPTX 생성 슬라이드 수 정확성 (Validates: Requirements 5.1, 5.3, 5.4).

For any non-empty slide list passed to _tool_generate_pptx:
  - response slideCount must equal len(slides) + 1 (cover slide)
  - the saved .pptx file must exist and be parseable by python-pptx
  - the parsed pptx slide count must match the reported slideCount
  - even when image generation fails for slides that carry an imagePrompt,
    slideCount stays the same (no slides dropped — Req 5.3)

Image generation is patched with unittest.mock so the test never hits Bedrock;
hypothesis flips a per-slide bit to decide whether the patched
_tool_generate_image returns a real PNG path (success) or an error JSON (fail),
exercising the fail-soft branch in _tool_generate_pptx.

Run:
  ai_engine/.venv/bin/python scripts/test_pptx_slide_count_property.py
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile
from io import BytesIO
from unittest.mock import patch

# Make the ai_engine package importable from the repo root.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402

from ai_engine.server import _tool_generate_pptx  # noqa: E402


@pytest.fixture(autouse=True)
def _disable_native_inference(monkeypatch):
    """이 파일은 imagePrompt→이미지 임베드 슬라이드수 불변식을 검증한다.
    네이티브 다이어그램 추론이 켜져 있으면 'TREE' 등 다이어그램형 본문이
    네이티브 도형으로 라우팅돼 _tool_generate_image 가 호출되지 않으므로
    (의도된 동작), 본 테스트 한정으로 추론을 꺼 이미지 경로를 순수 검증한다.
    """
    monkeypatch.setenv("AE_DISABLE_NATIVE_DIAGRAM", "1")


# ---------- helpers ----------

def _png_bytes(w: int = 8, h: int = 8) -> bytes:
    buf = BytesIO()
    Image.new("RGB", (w, h), (200, 200, 200)).save(buf, format="PNG")
    return buf.getvalue()


_PNG_BYTES = _png_bytes()


class _ImageFake:
    """Patched replacement for _tool_generate_image.

    Each call consumes the next outcome flag from `outcomes` (True=success,
    False=failure). On success, writes a real PNG into .generated/ inside
    `project_path` so _tool_generate_pptx's `os.path.isfile` check passes
    and add_picture can embed it. On failure, returns an error JSON exactly
    like the production fallback-exhausted path would.
    """

    def __init__(self, outcomes: list[bool]):
        self.outcomes = list(outcomes)
        self.idx = 0
        self.calls = 0

    async def __call__(self, tool_input, project_path, aws_profile="", bedrock_user=""):
        i = self.idx
        self.idx += 1
        self.calls += 1
        # Default to failure if hypothesis under-supplied outcomes (defensive).
        ok = self.outcomes[i] if i < len(self.outcomes) else False
        if not ok:
            return json.dumps({"error": "model-unavailable", "detail": "patched failure"})

        gen_dir = os.path.join(project_path, ".generated")
        os.makedirs(gen_dir, exist_ok=True)
        filename = f"fake-image-{i}.png"
        abs_path = os.path.join(gen_dir, filename)
        with open(abs_path, "wb") as f:
            f.write(_PNG_BYTES)
        return json.dumps({
            "path": f".generated/{filename}",
            "model": "patched.fake-image-model",
            "size": "1024x1024",
        })


# ---------- hypothesis strategies ----------

# Slide titles/bodies use printable ASCII (avoid pptx XML escaping edge cases
# noise — the property is about slide count, not title rendering).
_safe_text = st.text(
    alphabet=st.characters(min_codepoint=0x20, max_codepoint=0x7E),
    min_size=0,
    max_size=40,
)

_layout = st.sampled_from(["title", "content", "two-column"])

# Each slide: optional imagePrompt drives whether the patched generator is
# invoked at all, and `image_ok` decides success/fail when it is.
_slide_strategy = st.fixed_dictionaries({
    "title": _safe_text,
    "body": _safe_text,
    "layout": _layout,
    "has_image": st.booleans(),
    "image_ok": st.booleans(),
})


def _build_input(slide_specs):
    slides = []
    expected_image_calls = 0
    expected_image_outcomes: list[bool] = []
    for i, spec in enumerate(slide_specs):
        slide = {
            "title": spec["title"] or f"Slide {i + 1}",
            "layout": spec["layout"],
            "bullets": [spec["body"][:30]] if spec["body"] else [],
        }
        if spec["has_image"]:
            slide["imagePrompt"] = f"prompt for slide {i}"
            expected_image_calls += 1
            expected_image_outcomes.append(spec["image_ok"])
        slides.append(slide)
    return slides, expected_image_calls, expected_image_outcomes


# ---------- the property ----------

@settings(
    max_examples=40,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow, HealthCheck.function_scoped_fixture],
)
@given(slide_specs=st.lists(_slide_strategy, min_size=1, max_size=15))
def test_pptx_slide_count_invariant(slide_specs):
    slides, _expected_calls, image_outcomes = _build_input(slide_specs)

    fake_image = _ImageFake(image_outcomes)

    with tempfile.TemporaryDirectory() as tmp:
        tool_input = {
            "title": "Property Test Deck",
            "slides": slides,
        }

        with patch("ai_engine.server._tool_generate_image", new=fake_image):
            raw = asyncio.run(_tool_generate_pptx(tool_input, project_path=tmp))

        parsed = json.loads(raw)

        # 1. No error path — slides list is non-empty, so generation must succeed.
        assert "error" not in parsed, f"unexpected error: {parsed}"

        # 2. Req 5.4 — slideCount = len(slides) + 1 (cover included).
        expected_count = len(slides) + 1
        assert parsed.get("slideCount") == expected_count, (
            f"slideCount mismatch: got {parsed.get('slideCount')}, "
            f"expected {expected_count}; parsed={parsed}"
        )

        # 3. File exists on disk under the temp project root.
        rel = parsed.get("path", "")
        assert rel.startswith(".generated/") and rel.endswith(".pptx"), (
            f"unexpected path: {rel}"
        )
        abs_path = os.path.join(tmp, rel)
        assert os.path.isfile(abs_path), f"pptx not saved at {abs_path}"
        assert os.path.getsize(abs_path) > 0, "pptx file is empty"

        # 4. python-pptx can parse it AND the parsed slide count matches the
        #    reported slideCount. This catches any silent slide-drop, which
        #    is exactly what Req 5.3 forbids when image generation fails.
        prs = Presentation(abs_path)
        assert len(prs.slides) == expected_count, (
            f"parsed slide count {len(prs.slides)} != reported {expected_count} "
            f"(image_outcomes={image_outcomes})"
        )

        # 5. Patched _tool_generate_image was invoked exactly once per slide
        #    that declared an imagePrompt — proves slides aren't being skipped
        #    before the image step either.
        assert fake_image.calls == sum(1 for s in slides if "imagePrompt" in s), (
            f"image-generation call count mismatch: "
            f"calls={fake_image.calls}, "
            f"expected={sum(1 for s in slides if 'imagePrompt' in s)}"
        )


# ---------- focused sanity case: every imagePrompt fails ----------

def test_all_image_failures_preserve_slide_count():
    """Sanity case for Req 5.3 — every image fails, slide count must hold.

    Hypothesis already exercises this stochastically, but a deterministic
    smoke check keeps regressions obvious in CI logs.
    """
    n = 5
    slides = [
        {
            "title": f"Slide {i}",
            "layout": "content",
            "bullets": [f"point {i}"],
            "imagePrompt": f"prompt {i}",
        }
        for i in range(n)
    ]
    fake_image = _ImageFake([False] * n)

    with tempfile.TemporaryDirectory() as tmp:
        with patch("ai_engine.server._tool_generate_image", new=fake_image):
            raw = asyncio.run(_tool_generate_pptx(
                {"title": "All Fail Deck", "slides": slides},
                project_path=tmp,
            ))
        parsed = json.loads(raw)

    assert parsed.get("slideCount") == n + 1, parsed
    assert "error" not in parsed, parsed
    assert fake_image.calls == n, fake_image.calls


def main():
    print("=== Property 8: PPTX 슬라이드 수 정확성 ===")
    test_all_image_failures_preserve_slide_count()
    print("  deterministic all-fail case                     OK")
    test_pptx_slide_count_invariant()
    print("  hypothesis property (slideCount invariant)      OK")
    print("All Property 8 cases passed.")


if __name__ == "__main__":
    main()
