#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""KPI/진행률 네이티브 아키타입 회귀 테스트 (젠스파크 OKR 스타일).

목적: 편집 가능한 도형으로 KPI 요약 카드 + 진행률 막대를 렌더하고,
_classify_section_diagram이 OKR/지표/진척형 콘텐츠를 올바르게 라우팅하는지 검증.
편집 가능성(통짜 PICTURE 금지)을 핵심 불변식으로 검사한다.
"""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ai_engine.native_diagram_pptx import (
    build_native_diagram,
    _parse_kpis,
    _parse_progress,
)


def _blank_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _count_shapes(slide):
    auto = pic = 0
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pic += 1
        elif sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            auto += 1
    return auto, pic


# ---------- 파서 ----------

def test_parse_kpis_formats():
    kpis = _parse_kpis(
        "매출 증가: 32% (+8%p)\n"
        "고객 만족도 | 95% | +3%\n"
        "신규 고객 320명\n"
        "재구매율: 68% (-2%)"
    )
    assert len(kpis) == 4
    assert kpis[0][0] == "32%"
    assert kpis[0][1] == "매출 증가"
    assert kpis[0][2].startswith("+8")
    assert kpis[1][0] == "95%" and kpis[1][2].startswith("+3")
    assert kpis[2][0].startswith("320")
    assert kpis[3][2].startswith("-2")


def test_parse_kpis_skips_non_numeric():
    kpis = _parse_kpis("설명만 있는 줄\n또 다른 텍스트")
    assert kpis == []


def test_parse_progress_formats():
    prog = _parse_progress(
        "목표 A 시장 점유율 확대: 80% 진행중\n"
        "목표 B 제품 출시: 100% 완료\n"
        "목표 C 글로벌 진출: 45% 지연\n"
        "퍼센트 없는 설명 줄"
    )
    assert len(prog) == 3
    assert prog[0][1] == 80 and prog[0][2] == "진행중"
    assert prog[1][1] == 100 and prog[1][2] == "완료"
    assert prog[2][1] == 45 and prog[2][2] == "지연"


def test_parse_progress_clamps_percent():
    prog = _parse_progress("초과 항목: 150%\n음수는 패턴상 불가")
    assert prog[0][1] == 100  # clamp 0~100


# ---------- 렌더 (편집 가능 도형, PICTURE 0개) ----------

def test_kpi_renders_editable_shapes():
    prs, slide = _blank_slide()
    ok = build_native_diagram(
        slide, "kpi",
        "매출 증가: 32% (+8%p)\n고객 만족도 | 95% | +3%\n신규 고객 320명\n재구매율: 68%",
        title="핵심 성과 지표",
    )
    auto, pic = _count_shapes(slide)
    assert ok is True
    assert pic == 0, "KPI 슬라이드에 통짜 PICTURE가 있으면 안 됨(편집 가능 불변식)"
    assert auto >= 4, f"카드 도형 부족: {auto}"


def test_progress_renders_editable_shapes():
    prs, slide = _blank_slide()
    ok = build_native_diagram(
        slide, "progress",
        "목표 A 시장 점유율: 80% 진행중\n목표 B 제품 출시: 100% 완료\n목표 C 글로벌: 45% 지연",
        title="OKR 진척 현황",
    )
    auto, pic = _count_shapes(slide)
    assert ok is True
    assert pic == 0, "진행률 슬라이드에 통짜 PICTURE가 있으면 안 됨"
    # 행마다 최소 트랙+채움 도형 → 3행이면 6개 이상
    assert auto >= 6, f"진행 막대 도형 부족: {auto}"


def test_kpi_empty_content_returns_false():
    prs, slide = _blank_slide()
    assert build_native_diagram(slide, "kpi", "") is False
    assert build_native_diagram(slide, "kpi", "숫자 없는 텍스트만") is False


def test_progress_empty_content_returns_false():
    prs, slide = _blank_slide()
    assert build_native_diagram(slide, "progress", "") is False
    assert build_native_diagram(slide, "progress", "퍼센트 없는 텍스트") is False


# ---------- 분류기 라우팅 ----------

def test_classify_routes_kpi():
    from ai_engine.server import _classify_section_diagram
    kind, _ = _classify_section_diagram(
        "핵심 성과 지표", "매출 32%, 고객 만족도 95%, 신규 320명"
    )
    assert kind == "kpi"


def test_classify_routes_progress():
    from ai_engine.server import _classify_section_diagram
    kind, _ = _classify_section_diagram(
        "OKR 진척 현황", "목표 A: 80% 진행중\n목표 B: 100% 완료\n목표 C: 45%"
    )
    assert kind == "progress"


def test_classify_preserves_existing_kinds():
    from ai_engine.server import _classify_section_diagram
    assert _classify_section_diagram("프로젝트 흐름", "계획 -> 개발 -> 출시")[0] == "flow"
    assert _classify_section_diagram("디렉토리 구조", "  src/\n  tests/\n  docs/")[0] == "tree"


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v", "-p", "no:cacheprovider"]))
