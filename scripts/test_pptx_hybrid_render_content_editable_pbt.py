# Feature: pptx-ultra-quality-hybrid-render, Property 7: content 슬라이드는 항상 편집 가능
#   (For any content 슬라이드를 헤르메틱하게(Vertex 비활성) 렌더한 산출물에 대해,
#    html_enabled 값과 무관하게 편집 가능 텍스트 run 개수 >= 1 이고 슬라이드 전체
#    (13.333in x 7.5in)를 덮는 풀블리드 PICTURE 개수 == 0 이다.  Validates R2.1, R2.2)
# Feature: pptx-ultra-quality-hybrid-render, Property 8: content 히어로의 바운디드 합성·보존
#   (For any content 슬라이드와 유효 히어로 rel에서 보존/합성 이미지 개수 >= 1 이고
#    어떤 이미지도 풀블리드 PICTURE가 아니다; 슬롯 미지원 레이아웃에서도 미폐기.
#    Validates R2.3, R2.4)
# Feature: pptx-ultra-quality-hybrid-render, Property 9: content 산출물 밀도·스타일 감사 통과
#   (For any content 슬라이드를 포함한 헤르메틱 렌더 산출물과 tokens에 대해
#    audit_native_density(pptx_path, tokens)는 AuditReport.passed == True 이고
#    failures == [] 이다.  Validates R2.5, R5.3)
#
# 실행 규약(design.md §Testing Strategy / tasks.md 필수 실행 규약):
#   - 헤르메틱: Vertex 비활성(_render_content_editable 은 Vertex/네트워크를 호출하지 않는
#     순수 조립 함수이며, hero_rel 은 로컬 파일 경로다) — 외부 네트워크 호출 0.
#   - 각 property 테스트 최소 100 iteration(hypothesis max_examples=100).
#   - 실행: ./venv/bin/python -m pytest scripts/test_pptx_hybrid_render_content_editable_pbt.py \
#           -p no:cacheprovider -q
#   - 인라인 멀티라인 python -c 미사용(재현 가능한 스크립트 파일).
"""Property-based tests (Property 7/8/9) for ``ai_engine.server._render_content_editable``.

대상은 이미 구현된 조립 함수
``_render_content_editable(slide, prs, data, tokens, hero_rel, palette)`` 이다(task 5.2).
본 테스트는 구현 코드를 수정하지 않고, 실제 python-pptx 산출물을 열어 측정 가능한
불변식(편집 가능 텍스트 run, 풀블리드 PICTURE 부재, 이미지 바운디드 합성/보존, 밀도·스타일
감사 통과)을 검증한다.
"""

import os
import string
import sys
import tempfile

from hypothesis import given, settings, strategies as st, HealthCheck

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

# --- import 경로 정비: repo root + scripts/ ---------------------------------
_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _ROOT not in sys.path:
    sys.path.insert(0, _ROOT)
_SCRIPTS = os.path.join(_ROOT, "scripts")
if _SCRIPTS not in sys.path:
    sys.path.insert(0, _SCRIPTS)

import ai_engine.server as srv  # noqa: E402  (대상 함수 보유 모듈)
import ai_engine.native_layout_renderer as nlr  # noqa: E402  (design_tokens 재사용)
import audit_pptx_native_density as auditor  # noqa: E402  (Property 9 감사기)

from PIL import Image  # noqa: E402


# ===========================================================================
# 상수 — 풀블리드(슬라이드 전체) 판정 (16:9, 13.333in x 7.5in)
# ===========================================================================
_FULL_W_EMU = 12192000   # 13.333in
_FULL_H_EMU = 6858000    # 7.5in
_FB_ORIGIN_TOL = 6000    # 좌상단 원점 허용 오차(EMU, ~0.0007in)
_FB_SIZE_TOL = 6000      # 전체 크기 허용 오차(EMU)


# ===========================================================================
# 헤르메틱 성립 확인 — _render_content_editable 은 Vertex/네트워크를 호출하지 않는다.
# get_vertex_image_client 을 disabled 스텁으로 대체해 "Vertex 비활성"을 명시적으로
# 고정한다(어떤 우회 경로도 실 클라이언트를 취득하지 못하게).
# ===========================================================================
class _DisabledVertexClient:
    enabled = False

    def resolve_model_id(self, *_a, **_k):
        return ""

    async def generate(self, *_a, **_k):  # pragma: no cover - 호출되지 않아야 함
        raise AssertionError("헤르메틱 위반: Vertex generate 가 호출됨")


try:  # 대상 함수는 Vertex 를 만지지 않지만, 안전하게 disabled 로 고정한다.
    import ai_engine.vertex_image_module as _vim  # noqa: E402

    _vim.get_vertex_image_client = lambda *a, **k: _DisabledVertexClient()  # type: ignore[assignment]
except Exception:
    pass


_COMMON = settings(
    max_examples=100,          # 각 property 최소 100 iteration
    deadline=None,             # 실제 .pptx I/O 비용 반영
    suppress_health_check=[HealthCheck.too_slow, HealthCheck.data_too_large,
                           HealthCheck.filter_too_much],
)


# ===========================================================================
# 생성기 — 한글/영문/숫자/특수문자 혼합, python-pptx XML 에 안전한 문자만.
# ===========================================================================
_SAFE_ALPHABET = (
    "가나다라마바사아자차카타파하거너더러머버서어전정종주중지매출성장고객"
    + string.ascii_letters
    + string.digits
    + " .,!?·—…:/@#%&*()-_=+"
)


def _ne_text(min_size=1, max_size=40):
    return (st.text(alphabet=_SAFE_ALPHABET, min_size=min_size, max_size=max_size)
            .map(lambda s: s.strip())
            .filter(lambda s: len(s) > 0))


# Property 7/8: 넓은 content 데이터(불릿 0~6개 — 리치/최소 레이아웃 모두 커버).
_bullets_wide = st.lists(_ne_text(1, 30), min_size=0, max_size=6)

# Property 9: 밀도 감사 통과를 위해 content 슬라이드는 실제 밀도(불릿 >= 2)를 가진다.
_bullets_dense = st.lists(_ne_text(1, 30), min_size=2, max_size=6)

# Style_Profile — normalize 대상 색상(대소문자 혼합) + None 폴백.
_HEX = st.text(alphabet="0123456789abcdefABCDEF", min_size=6, max_size=6)


@st.composite
def _profile(draw):
    if draw(st.booleans()):
        return None
    return {
        "primaryColor": "#" + draw(_HEX),
        "secondaryColor": "#" + draw(_HEX),
        "accentColor": "#" + draw(_HEX),
        "textColor": "#" + draw(_HEX),
    }


# Property 9 전용 — 현실적 Style_Profile 팔레트.
#   audit_native_density 의 (h) 스타일 품질 검사와 색 기반 시각 요소 검출기는 설계상
#   "스타일된" 색을 요구한다: 강조/주색 계열은 유채색(검정/흰색·근접 무채색이 아님),
#   본문 텍스트색은 어두운 색이다. 실제 Style_Profile(브랜드 팔레트)의 의미 구조가
#   정확히 이렇다 — primary/secondary/accent 는 선명한 유채색, textColor 는 짙은 색.
#   순수 흑/백·근접 무채색을 강조색 슬롯에 넣는 것은 어떤 디자인 프로파일도 하지 않는
#   퇴화 입력이므로 Property 9 의 입력 공간(=Style_Profile 도메인)에서 제외한다.
#   (_render_content_editable 은 주어진 tokens 를 충실히 적용하며, 색 무관 속성인
#   Property 7/8 은 넓은 팔레트로 이미 검증한다.)
_VIVID_HEX = [
    "#3B82C4", "#2E8B57", "#F5A623", "#0066FF", "#00C896", "#FF6B35",
    "#AB12CD", "#12CD34", "#C0392B", "#8E44AD", "#16A085", "#D35400",
    "#2980B9", "#27AE60", "#E67E22", "#6C5CE7", "#00B894",
]
_DARK_TEXT_HEX = ["#1A2332", "#111111", "#222222", "#1B1B2F", "#0F1724", "#141414"]


@st.composite
def _realistic_profile(draw):
    if draw(st.booleans()):
        return None
    cols = draw(st.lists(st.sampled_from(_VIVID_HEX),
                         min_size=3, max_size=3, unique=True))
    return {
        "primaryColor": cols[0],
        "secondaryColor": cols[1],
        "accentColor": cols[2],
        "textColor": draw(st.sampled_from(_DARK_TEXT_HEX)),
    }


# ===========================================================================
# python-pptx 헬퍼
# ===========================================================================
def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def _blank(prs):
    # layout[6] = Blank (add_slide 자동 placeholder 없는 빈 레이아웃).
    return prs.slides.add_slide(prs.slide_layouts[6])


def _editable_run_count(slide) -> int:
    """편집 가능 텍스트 프레임의 비어있지 않은 run 개수."""
    n = 0
    for sp in slide.shapes:
        try:
            if not getattr(sp, "has_text_frame", False):
                continue
            for para in sp.text_frame.paragraphs:
                for run in para.runs:
                    if (run.text or "").strip():
                        n += 1
        except Exception:
            continue
    return n


def _is_fullbleed_picture(sp) -> bool:
    """셰이프가 슬라이드 전체(0,0)~(13.333in x 7.5in)를 덮는 PICTURE 인지."""
    try:
        if sp.shape_type != MSO_SHAPE_TYPE.PICTURE:
            return False
        l, t, w, h = sp.left, sp.top, sp.width, sp.height
        if None in (l, t, w, h):
            return False
        return (abs(int(l)) <= _FB_ORIGIN_TOL and abs(int(t)) <= _FB_ORIGIN_TOL
                and int(w) >= _FULL_W_EMU - _FB_SIZE_TOL
                and int(h) >= _FULL_H_EMU - _FB_SIZE_TOL)
    except Exception:
        return False


def _count_pictures(slide):
    total = 0
    fullbleed = 0
    for sp in slide.shapes:
        try:
            if sp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                total += 1
                if _is_fullbleed_picture(sp):
                    fullbleed += 1
        except Exception:
            continue
    return total, fullbleed


def _write_small_png(path, size=(320, 240)):
    Image.new("RGB", size, (90, 130, 190)).save(path)


def _content_data(title, bullets, layout_hint=None):
    d = {"title": title, "bullets": list(bullets)}
    if layout_hint is not None:
        d["layout"] = layout_hint
    return d


# ===========================================================================
# Property 7 — content 슬라이드는 항상 편집 가능
# ===========================================================================
@_COMMON
@given(
    title=_ne_text(1, 48),
    bullets=_bullets_wide,
    profile=_profile(),
    html_enabled=st.booleans(),
)
def test_property7_content_always_editable(title, bullets, profile, html_enabled):
    """헤르메틱(Vertex 비활성) content 렌더는 html_enabled 값과 무관하게
    편집 가능 텍스트 run >= 1 이고 풀블리드 PICTURE 개수 == 0 이다.

    **Validates: Requirements 2.1, 2.2**
    """
    # html_enabled 독립성 명시: 전역 HTML 플래그를 켜고/꺼도 결과 불변이어야 함.
    prev = os.environ.get("AE_ENABLE_HTML_SLIDES")
    os.environ["AE_ENABLE_HTML_SLIDES"] = "1" if html_enabled else "0"
    try:
        tokens = nlr.design_tokens_for_profile(profile)
        prs = _new_prs()
        slide = _blank(prs)
        res = srv._render_content_editable(
            slide, prs, _content_data(title, bullets), tokens, None, None)

        # 조립 함수는 항상 편집 의도(editable=True)를 보고한다.
        assert res.get("editable") is True

        # 산출물 실측: 편집 가능 텍스트 run >= 1 (R2.1/R2.2).
        runs = _editable_run_count(slide)
        assert runs >= 1, f"편집 가능 텍스트 run 0개 (html_enabled={html_enabled}, res={res})"

        # 산출물 실측: 슬라이드 전체를 덮는 풀블리드 PICTURE 0개 (R2.1).
        _total, fb = _count_pictures(slide)
        assert fb == 0, f"풀블리드 PICTURE {fb}개 발견 (편집성 위반)"
    finally:
        if prev is None:
            os.environ.pop("AE_ENABLE_HTML_SLIDES", None)
        else:
            os.environ["AE_ENABLE_HTML_SLIDES"] = prev


# ===========================================================================
# Property 8 — content 히어로의 바운디드 합성·보존
# ===========================================================================
@_COMMON
@given(
    title=_ne_text(1, 48),
    # 슬롯 지원(불릿 >=2 -> feature_grid) 과 슬롯 미지원(불릿 0 -> section_divider) 모두 커버.
    bullets=st.one_of(st.just([]), st.lists(_ne_text(1, 30), min_size=2, max_size=6)),
    profile=_profile(),
)
def test_property8_hero_bounded_composite_or_preserve(title, bullets, profile):
    """유효 히어로 rel 에서 보존/합성 이미지 개수 >= 1 이고 어떤 이미지도 풀블리드
    PICTURE 가 아니다(바운디드 슬롯 또는 바운디드 on-slide 레이어). 슬롯 미지원
    레이아웃에서도 이미지는 폐기되지 않는다.

    **Validates: Requirements 2.3, 2.4**
    """
    tokens = nlr.design_tokens_for_profile(profile)
    fd, hero = tempfile.mkstemp(suffix=".png", prefix="pbt_hero_")
    os.close(fd)
    _write_small_png(hero)
    try:
        prs = _new_prs()
        slide = _blank(prs)
        res = srv._render_content_editable(
            slide, prs, _content_data(title, bullets), tokens, hero, None)

        # 합성(슬롯) 또는 보존(레이어) 중 하나는 반드시 발생(폐기 금지, 손실-0/R2.4).
        assert res.get("image_placed") or res.get("image_preserved"), (
            f"히어로가 합성·보존 어느 쪽으로도 배치되지 않음 (폐기 의심, res={res})")

        # 산출물 실측: 이미지(PICTURE) 개수 >= 1.
        total, fb = _count_pictures(slide)
        assert total >= 1, f"보존/합성 이미지 0개 (res={res})"

        # 어떤 이미지도 풀블리드가 아님(바운디드, R2.3/R2.4).
        assert fb == 0, f"히어로가 풀블리드 PICTURE 로 배치됨 {fb}개 (바운디드 위반)"

        # 편집성도 함께 유지(히어로가 텍스트를 밀어내지 않음).
        assert _editable_run_count(slide) >= 1
    finally:
        try:
            os.remove(hero)
        except OSError:
            pass


# ===========================================================================
# Property 9 — content 산출물 밀도·스타일 감사 통과
# ===========================================================================
@_COMMON
@given(
    title=_ne_text(2, 48),
    bullets=_bullets_dense,
    profile=_realistic_profile(),
)
def test_property9_content_density_style_audit_passes(title, bullets, profile):
    """헤르메틱 content 렌더 산출물과 tokens 에 대해 audit_native_density 가
    AuditReport.passed == True 이고 failures == [] 이다(비텍스트 시각 요소 >= 2 및
    5개 스타일 품질 검사 통과 포함).

    감사기는 슬라이드 인덱스 0 을 cover 카테고리로, 이후를 body 카테고리로 채점한다.
    content 슬라이드는 실전에서 표지 다음(body)에 온다 — 따라서 표지(cover) 1장을
    앞세운 뒤 content 슬라이드를 두 번째로 렌더해 body 카테고리로 감사한다.

    **Validates: Requirements 2.5, 5.3**
    """
    tokens = nlr.design_tokens_for_profile(profile)
    prs = _new_prs()

    # 실제 content 는 제목과 불릿 텍스트가 서로 구별된다. 감사기의 (d) 제목-유일성
    # 검사는 동일 문자열이 여러 셰이프에 나타나면 "제목 2개"로 본다(제목이 불릿과
    # 우연히 같은 텍스트인 퇴화 입력). 라벨 접두사로 제목≠불릿, 불릿 상호 구별을
    # 보장한다(밀도/스타일 검사에는 영향 없음 — 텍스트 내용만 라벨링).
    title_t = "제목· " + title
    bullets_t = [f"{i + 1}· {b}" for i, b in enumerate(bullets)]

    # 슬라이드 0: cover (감사 cover 카테고리 충족 — content 슬라이드가 cover 로
    #             오분류되지 않도록 하는 스캐폴딩. render_native_layout 편집 경로).
    cover = _blank(prs)
    nlr.render_native_layout(
        cover, prs, "cover",
        {"title": "프레젠테이션 표지", "subtitle": "부제 텍스트",
         "eyebrow": "OVERVIEW", "footer": "2026"},
        tokens)

    # 슬라이드 1: content — 대상 함수(편집 경로). 이 슬라이드가 body 로 감사된다.
    content = _blank(prs)
    res = srv._render_content_editable(
        content, prs, _content_data(title_t, bullets_t), tokens, None, None)
    assert res.get("editable") is True

    fd, path = tempfile.mkstemp(suffix=".generated.pptx", prefix="pbt_content_density_")
    os.close(fd)
    try:
        prs.save(path)
        report = auditor.audit_native_density(path, tokens=tokens)
    finally:
        try:
            os.remove(path)
        except OSError:
            pass

    assert report.passed is True, (
        "content 산출물 밀도·스타일 감사 미통과:\n"
        + "\n".join(f"  - {f['check']} (slide {f['slide']}): {f['signal']}"
                    for f in report.failures))
    assert report.failures == []


if __name__ == "__main__":
    # 직접 실행 시 각 property 를 한 번씩 스모크(예시 입력)로 구동.
    test_property7_content_always_editable()
    test_property8_hero_bounded_composite_or_preserve()
    test_property9_content_density_style_audit_passes()
    print("OK: Property 7/8/9 smoke passed")
