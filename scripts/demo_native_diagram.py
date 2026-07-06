"""정교화된 네이티브(편집 가능) 다이어그램 데모.

flow/tree/block 3종을 각 슬라이드에 네이티브 도형으로 그려
.generated/native-diagram-demo.pptx 로 저장. PowerPoint/Keynote에서 열어 도형·텍스트·
화살표가 개별 편집되는지, 정렬/간격/번호배지/섀도우 품질을 확인.

실행: python scripts/demo_native_diagram.py
"""
from __future__ import annotations
import sys
from pathlib import Path

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT)); sys.path.insert(0, str(_ROOT / "ai_engine"))

from pptx import Presentation
from pptx.util import Inches, Pt
import native_diagram_pptx as nd

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
REGION = (0.7, 1.6, 11.9, 5.4)

def _slide(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = s.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(0.9))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True
    return s

# 1) 가로 흐름
s1 = _slide("① EC2 비용 모니터링 에이전트 흐름")
nd.build_native_diagram(
    s1, "flow",
    "GPU/CPU 모드 자동 감지 -> psutil 프로세스 스캔 -> PID 추적 -> 프로젝트 식별",
    region=REGION, palette=["#2E5BBA", "#F5A623"])

# 2) 세로 흐름(6+)
s2 = _slide("② MOGAM 멀티 에이전트 시스템 흐름")
nd.build_native_diagram(
    s2, "flow",
    "Team Lead 라우터 -> AI 기반 질문 분류 -> AI 에이전트 -> LNP 에이전트 -> Bioinfo 에이전트 -> Lead 통합",
    region=REGION, palette=["#2E5BBA", "#00897B"])

# 3) 조직도(트리)
s3 = _slide("③ 통합 아키텍처 구성")
nd.build_native_diagram(
    s3, "tree",
    "AI 비용·운영 플랫폼\n  비용 집계 ETL\n  통합 모니터링\n  계정 자동화\n  AI 사용량 거버넌스",
    region=REGION, palette=["#2E5BBA", "#5B9BD5"])

# 4) 블록(특징 나열)
s4 = _slide("④ 주요 구현 특징")
nd.build_native_diagram(
    s4, "block",
    "agent_new.py — EC2 인스턴스별 비용 추적 에이전트\n"
    "smart_agents_v2.py — 5개 Bedrock Agent 통합 관리\n"
    "실시간 스트리밍 — 60초 주기 메트릭 + 즉시 응답 Q&A\n"
    "프로젝트 식별 — Git repo > FSX 패턴 > cmdline 우선순위",
    region=REGION, palette=["#F5A623", "#2E5BBA"])

# 5) 카드 그리드(젠스파크풍, 편집 가능, 비통짜)
s5 = _slide("⑤ 프로젝트 개요 — 카드 그리드")
nd.build_native_diagram(
    s5, "cards",
    "체계적인 구조: 총 93개 디렉토리로 구성된 체계적인 프로젝트 구조\n"
    "효율적인 분류: 5가지 카테고리로 자동 분류된 폴더 관리\n"
    "확장 가능한 아키텍처: 멀티 레벨 아키텍처 기반 확장 가능 환경\n"
    "DevOps 중심: DevOps 중심의 현대적 프로젝트 관리 접근법\n"
    "모듈화된 설계: 유지보수성과 재사용성 극대화",
    region=REGION, palette=["#2E5BBA", "#00897B"])

gen = _ROOT / ".generated"; gen.mkdir(parents=True, exist_ok=True)
out = gen / "native-diagram-demo.pptx"
prs.save(str(out))
print("saved:", out)
print("PowerPoint/Keynote로 열어 도형·텍스트·화살표 개별 편집 + 정렬/간격/번호배지/섀도우를 확인하세요.")
