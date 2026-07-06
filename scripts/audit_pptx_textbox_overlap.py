"""오버레이 텍스트 박스끼리의 겹침(text-on-text) 정밀 계산."""
from __future__ import annotations
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU = 914400.0
def _in(v): return round(float(v)/EMU, 2) if v is not None else None

def boxes(slide):
    out = []
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            continue
        if getattr(sh, "has_text_frame", False) and (sh.text_frame.text or "").strip():
            l, t, w, h = _in(sh.left), _in(sh.top), _in(sh.width), _in(sh.height)
            if None in (l, t, w, h):
                continue
            head = (sh.text_frame.text or "").strip().split("\n")[0][:24]
            out.append((head, l, t, w, h))
    return out

def ov(a, b):
    _, l1, t1, w1, h1 = a
    _, l2, t2, w2, h2 = b
    ix = max(0.0, min(l1+w1, l2+w2) - max(l1, l2))
    iy = max(0.0, min(t1+h1, t2+h2) - max(t1, t2))
    return ix*iy

def main(path):
    prs = Presentation(path)
    print(f"=== 텍스트 박스끼리 겹침: {path}\n")
    for idx, slide in enumerate(prs.slides, 1):
        bs = boxes(slide)
        found = False
        for i in range(len(bs)):
            for j in range(i+1, len(bs)):
                a = ov(bs[i], bs[j])
                if a > 0.05:  # 0.05 in^2 이상 겹침
                    amin = min(bs[i][3]*bs[i][4], bs[j][3]*bs[j][4])
                    pct = a/amin*100 if amin else 0
                    print(f"[슬라이드 {idx}] '{bs[i][0]}' ↔ '{bs[j][0]}' "
                          f"겹침 {a:.2f}in² (작은박스의 {pct:.0f}%)")
                    found = True
        if not found:
            print(f"[슬라이드 {idx}] 텍스트 박스 간 겹침 없음")
    print()

if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "")
