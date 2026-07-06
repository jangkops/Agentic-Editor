"""Preservation property tests (P1–P4) — spec: media-output-quality (bugfix), Task 8.

These property-based tests validate the behaviours the fix (tasks 2–6, already
present in ``ai_engine/server.py``) MUST PRESERVE. Following the observation-first
methodology, each assertion was derived by first observing what the *current*
code actually does at the relevant seam, then asserting exactly that.

  P1 (Req 3.3 — circuit trip): a round where every Bedrock image model returns
     access-denied → ``_tool_generate_image`` trips the circuit
     (``_IMAGE_GEN_CIRCUIT["disabled_at"] != 0``).
  P2 (Req 3.4 — non-visual routing): a NON-visual-intent PPTX/PDF request →
     ``gw.invoke_model`` is never called with an image model id
     (``stability.*`` / ``amazon.nova-canvas-v1:0`` /
     ``amazon.titan-image-generator-v2:0``).
  P3 (Req 3.5 — PPTX layout coords): a slide with BOTH text and image → the body
     placeholder is positioned at ``(left=0.6, top=1.6, width=6.0, height=5.4)``
     and the image is positioned at ``x == Inches(7.0)``.
  P4 (Req 3.7 — TTL auto-recovery): a force-tripped circuit whose TTL has elapsed
     auto-recovers (``_image_gen_is_circuit_broken() == False``) and the next
     ``_tool_generate_image`` call with a HEALTHY gateway re-fires image-model
     attempts.

OBSERVATION NOTE for P3 (design-literal reconciliation):
The two-column image seam lives in ``_tool_generate_pptx`` (the ``elif img_path:``
branch). Observing the current code:

    if _has_text:
        body_shape.left   = Emu(int(0.6 * 914400))   # Inches(0.6)
        body_shape.top    = Emu(int(1.6 * 914400))   # Inches(1.6)
        body_shape.width  = Emu(int(6.0 * 914400))   # Inches(6.0)
        body_shape.height = Emu(int(5.4 * 914400))   # Inches(5.4)
        region_l, region_t, region_w, region_h = 7.0, 1.6, 6.0, 5.4
    ...
    draw_w = region_w                       # 6.0
    draw_h = region_w / ar                  # 6.0 / aspect-ratio
    if draw_h > region_h: draw_h, draw_w = region_h, region_h * ar
    off_l = region_l + (region_w - draw_w) / 2.0     # image x (centered in region)

* The BODY placeholder coordinates match the design literals EXACTLY
  (0.6 / 1.6 / 6.0 / 5.4).
* The IMAGE x-coordinate equals the region's left edge (``region_l == 7.0``) ONLY
  when the picture fills the full region width (``draw_w == region_w == 6.0``),
  i.e. when the source image is wide enough (aspect ratio ≥ 6.0/5.4 ≈ 1.111). For
  square/portrait images the code CENTERS the picture inside the [7.0, 13.0]
  region, so ``x > 7.0``. This is the actual, preserved behaviour. To make the
  design literal ``x == Inches(7.0)`` deterministic, P3 drives the seam with a
  KNOWN WIDE PNG (aspect ratio ≥ 1.3) so the picture fills the region width and
  ``off_l`` collapses to exactly ``region_l == 7.0``. The design literals and the
  observed behaviour therefore AGREE for the wide-image case P3 exercises.

Run (hermetic — no network, gateway mocked):
  ./venv/bin/python -m pytest scripts/test_media_output_quality_preservation_pbt.py -p no:cacheprovider -q

_Requirements: 3.3, 3.4, 3.5, 3.7_
"""
from __future__ import annotations

import os
import io
import sys
import json
import time
import base64
import asyncio
import random
import tempfile
from unittest.mock import patch

# Make ai_engine importable from repo root.
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import pytest  # noqa: E402
from hypothesis import given, settings, strategies as st, assume, HealthCheck  # noqa: E402

import ai_engine.server as srv  # noqa: E402
from ai_engine.server import (  # noqa: E402
    _detect_visual_intent,
    _looks_structural,
    _image_gen_is_circuit_broken,
    _IMAGE_GEN_CIRCUIT,
    _IMAGE_GEN_ATTEMPTS,
    IMAGE_MODELS,
)

# Access-denied gateway error string — must match _IMAGE_GEN_DENY_PATTERNS so the
# circuit-trip path (all-denied) is reached.
_ACCESS_DENIED = (
    "AccessDeniedException: not authorized to perform execute-api:Invoke "
    "(HTTP 403)"
)


def _is_image_model_id(model_id: str) -> bool:
    """True iff model_id is one of the Bedrock image-generation models."""
    if not model_id:
        return False
    return (
        model_id.startswith("stability.")
        or model_id == "amazon.titan-image-generator-v2:0"
        or model_id == "amazon.nova-canvas-v1:0"
        or model_id in set(IMAGE_MODELS)
    )


async def _async_none(*args, **kwargs):
    """Stand-in for _try_vertex_image_single — never returns a Vertex image."""
    return None


class _DisabledVertexClient:
    """Vertex client stub reporting itself disabled (no network)."""
    enabled = False
    _project_id = None


# --------------------------------------------------------------------------
# Gateway fakes
# --------------------------------------------------------------------------
class _DeniedGW:
    """Every invoke_model returns an access-denied error (records model ids)."""

    def __init__(self):
        self.calls: list[str] = []

    async def invoke_model(self, model_id, body, timeout=60):
        self.calls.append(model_id)
        return {"error": _ACCESS_DENIED}


class _SpyGW:
    """Records every invoke_model model id; returns a benign error so nothing
    proceeds. Used to prove image models are never invoked (P2)."""

    def __init__(self):
        self.calls: list[str] = []

    async def invoke_model(self, model_id, body, timeout=60):
        self.calls.append(model_id)
        return {"error": "spy: disabled in test"}


class _HealthyGW:
    """Every invoke_model returns a valid PNG (records model ids)."""

    def __init__(self, png_b64: str):
        self._png = png_b64
        self.calls: list[str] = []

    async def invoke_model(self, model_id, body, timeout=60):
        self.calls.append(model_id)
        return {"images": [self._png]}


def _make_valid_png_b64(seed: int = 7) -> str:
    """base64 PNG > 5KB with real entropy so _save_and_score accepts it."""
    from PIL import Image
    rng = random.Random(seed)
    img = Image.new("RGB", (256, 256))
    img.putdata([
        (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
        for _ in range(256 * 256)
    ])
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return base64.b64encode(buf.getvalue()).decode("ascii")


_VALID_PNG_B64 = _make_valid_png_b64()


# ==========================================================================
# P1 (Req 3.3) — all image models access-denied → circuit trips
# ==========================================================================
_PROMPTS = [
    "프로젝트 아키텍처 다이어그램",
    "system architecture diagram",
    "데이터 파이프라인 흐름도",
    "marketing hero illustration",
    "회사 소개 인포그래픽",
]


@settings(
    max_examples=12,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow],
)
@given(prompt=st.sampled_from(_PROMPTS))
def test_p1_all_denied_trips_circuit(prompt):
    gw = _DeniedGW()
    _IMAGE_GEN_CIRCUIT["disabled_at"] = 0  # start healthy
    _IMAGE_GEN_ATTEMPTS.clear()
    saved = {
        k: os.environ.get(k)
        for k in ("AE_PREFER_VERTEX_IMAGE", "AE_IMAGE_QUALITY_THRESHOLD")
    }
    os.environ["AE_PREFER_VERTEX_IMAGE"] = "0"     # skip vertex-first path
    os.environ["AE_IMAGE_QUALITY_THRESHOLD"] = "0"  # no quality-retry recursion
    try:
        with tempfile.TemporaryDirectory() as tmp, \
                patch.object(srv, "_get_gw", return_value=gw), \
                patch.object(srv, "_resolve_callable_model_id",
                             side_effect=lambda m, *a, **k: m), \
                patch.object(srv, "_try_vertex_image_single", new=_async_none), \
                patch("ai_engine.vertex_image_module.get_vertex_image_client",
                      return_value=_DisabledVertexClient()):
            asyncio.run(srv._tool_generate_image(
                {"prompt": prompt, "size": "1024x1024"},
                project_path=tmp,
            ))
        # The circuit MUST trip when every Bedrock image model is denied.
        assert _IMAGE_GEN_CIRCUIT["disabled_at"] != 0, (
            "circuit should trip (disabled_at != 0) when all image models are "
            f"access-denied; got {_IMAGE_GEN_CIRCUIT['disabled_at']!r}"
        )
        # Sanity: the trip was reached through REAL image-model attempts, not a
        # short-circuit — so the assertion above is not vacuous.
        assert any(_is_image_model_id(m) for m in gw.calls), (
            f"expected image-model invocations before trip, got: {gw.calls!r}"
        )
    finally:
        for k, v in saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v
        _IMAGE_GEN_CIRCUIT["disabled_at"] = 0
        _IMAGE_GEN_ATTEMPTS.clear()


# ==========================================================================
# P2 (Req 3.4) — non-visual request → zero image-model invocations
# ==========================================================================
# Non-visual sentence pools — deliberately avoid EVERY keyword in
# _detect_visual_intent (KOR: 시각/이미지/그림/사진/다이어그램/차트/그래프/흐름/
# 프로세스/분석/아키텍처/시스템/구조/슬라이드/발표/보고/개요/단계 …, EN: visual/
# image/diagram/chart/flow/process/architecture/system/slide/presentation/
# overview/structure …). An `assume` guard below enforces this hermetically.
_NON_VISUAL_SENTENCES = [
    "이번 달 매출 합계는 전월보다 늘었습니다.",
    "지난주 코드에서 함수 이름을 바꾸고 오류를 고쳤습니다.",
    "고객 문의 응답 시간이 평균 두 시간 줄었습니다.",
    "예산 집행 내역과 잔액을 정리했습니다.",
    "팀 회의에서 다음 분기 목표를 정했습니다.",
    "The team fixed several bugs and merged the pull request.",
    "Monthly revenue increased compared to last month.",
    "We updated the changelog and tagged a new release.",
    "Customer response time was reduced by two hours.",
    "The quarterly budget was approved by the finance lead.",
]
_NON_VISUAL_TITLES = [
    "월간 매출 정리",
    "코드 변경 내역",
    "분기 예산 요약",
    "Monthly Revenue Note",
    "Code Change Log",
    "Quarterly Budget Memo",
]


@settings(
    max_examples=20,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow, HealthCheck.filter_too_much],
)
@given(
    title=st.sampled_from(_NON_VISUAL_TITLES),
    sentences=st.lists(st.sampled_from(_NON_VISUAL_SENTENCES),
                       min_size=2, max_size=5),
    ext=st.sampled_from(["pptx", "pdf"]),
)
def test_p2_non_visual_request_no_image_calls(title, sentences, ext):
    body = " ".join(sentences)
    description = f"{title} 문서를 만들어줘"

    # Hermetic guard: the request MUST be non-visual and non-structural so the
    # image tiers are never eligible. (Pools are curated to satisfy this; the
    # assume is a belt-and-suspenders safety net.)
    assume(not _detect_visual_intent(description))
    assume(not _detect_visual_intent(title))
    assume(not _detect_visual_intent(body))
    assume(not _looks_structural(description, title, body))

    gw = _SpyGW()
    env_off = {
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_DISABLE_HTML_SLIDES": "1",
        "AE_PPTX_TOC": "0",
        "AE_DISABLE_MERMAID": "1",
        "AE_ENABLE_BEDROCK_SLIDE_IMAGES": "0",
    }
    saved = {k: os.environ.get(k) for k in env_off}
    os.environ.update(env_off)
    try:
        with tempfile.TemporaryDirectory() as tmp, \
                patch.object(srv, "_get_gw", return_value=gw), \
                patch.object(srv, "_resolve_callable_model_id",
                             side_effect=lambda m, *a, **k: m), \
                patch.object(srv, "_try_vertex_image_single", new=_async_none), \
                patch("ai_engine.vertex_image_module.get_vertex_image_client",
                      return_value=_DisabledVertexClient()):
            asyncio.run(srv._force_generate_from_text(
                primary_tool=f"generate_{ext}",
                target_files=[f"out.{ext}"],
                title=title,
                description=description,
                final_text=body,
                project_path=tmp,
                aws_profile="test",
                bedrock_user="",
            ))
    finally:
        for k, v in saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v

    image_calls = [m for m in gw.calls if _is_image_model_id(m)]
    assert image_calls == [], (
        "non-visual request must not invoke any Bedrock image model, but these "
        f"image-model calls were made: {image_calls!r} (all calls: {gw.calls!r})"
    )


# ==========================================================================
# P3 (Req 3.5) — PPTX two-column layout coordinates are preserved
# ==========================================================================
def _write_wide_png(path: str, width: int, height: int, seed: int = 3) -> None:
    """Write a valid wide PNG (aspect ratio = width/height) to `path`."""
    from PIL import Image
    rng = random.Random(seed)
    img = Image.new("RGB", (width, height))
    img.putdata([
        (rng.randint(0, 255), rng.randint(0, 255), rng.randint(0, 255))
        for _ in range(width * height)
    ])
    img.save(path, format="PNG")


def _emu_inches(inches: float) -> int:
    """python-pptx Inches() in EMU — int(inches * 914400)."""
    return int(inches * 914400)


@settings(
    max_examples=12,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow],
)
@given(
    slide_title=st.text(
        alphabet="가나다라마바사아자차ABCDEFGHJKLMN 0123456789",
        min_size=1, max_size=24,
    ),
    bullets=st.lists(
        st.text(alphabet="가나다라마바사아자차ABCDEFG 0123456789", min_size=1, max_size=30),
        min_size=1, max_size=4,
    ),
    base_w=st.integers(min_value=900, max_value=1400),
    ratio_x10=st.integers(min_value=13, max_value=22),  # aspect ratio 1.3 .. 2.2
)
def test_p3_pptx_two_column_layout_coords(slide_title, bullets, base_w, ratio_x10):
    # Ensure bullets carry real text → _has_text True → two-column branch.
    bullets = [b for b in bullets if b.strip()]
    assume(bullets)
    assume(slide_title.strip())

    ratio = ratio_x10 / 10.0
    width = base_w
    height = max(1, round(width / ratio))  # aspect ratio >= 1.3 → fills region width

    from pptx import Presentation
    from pptx.util import Inches
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    # Disable every higher tier so the deterministic imageFile two-column seam
    # is the one exercised (no network, no LLM structuring, no Vertex, no TOC).
    env_off = {
        "AE_PREFER_EDITABLE_DIAGRAM": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_ENABLE_HTML_SLIDES": "0",
        "AE_DISABLE_HTML_SLIDES": "1",
        "AE_PPTX_TOC": "0",
        "AE_DISABLE_NATIVE_DIAGRAM": "1",
    }
    saved = {k: os.environ.get(k) for k in env_off}
    os.environ.update(env_off)
    try:
        with tempfile.TemporaryDirectory() as tmp:
            png_path = os.path.join(tmp, "wide.png")
            _write_wide_png(png_path, width, height)
            tool_input = {
                "title": "Deck",
                "slides": [{
                    "title": slide_title,
                    "bullets": bullets,
                    "imageFile": png_path,  # absolute → embedded without Bedrock
                }],
            }
            raw = asyncio.run(srv._tool_generate_pptx(tool_input, project_path=tmp))
            result = json.loads(raw)
            assert "error" not in result, f"pptx build failed: {result!r}"
            abs_path = result.get("absPath") or os.path.join(tmp, result["path"])
            assert os.path.isfile(abs_path), f"pptx not written: {abs_path}"

            prs = Presentation(abs_path)
            # slides[0] = cover; slides[1] = our content slide (TOC disabled).
            assert len(prs.slides) >= 2, "expected cover + content slide"
            content = prs.slides[1]

            # --- body placeholder coordinates (design literals, exact) ---
            body = content.placeholders[1]
            assert body.left == Inches(0.6), \
                f"body.left {body.left} != Inches(0.6) {Inches(0.6)}"
            assert body.top == Inches(1.6), \
                f"body.top {body.top} != Inches(1.6) {Inches(1.6)}"
            assert body.width == Inches(6.0), \
                f"body.width {body.width} != Inches(6.0) {Inches(6.0)}"
            assert body.height == Inches(5.4), \
                f"body.height {body.height} != Inches(5.4) {Inches(5.4)}"

            # --- image x-coordinate (design literal x == Inches(7.0)) ---
            pics = [sh for sh in content.shapes
                    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
            assert pics, "expected an embedded picture on the content slide"
            pic = pics[0]
            assert pic.left == Inches(7.0), (
                f"image x {pic.left} != Inches(7.0) {Inches(7.0)} "
                f"(wide image ar={width / height:.3f} should fill region width)"
            )
    finally:
        for k, v in saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v


# ==========================================================================
# P4 (Req 3.7) — TTL elapsed → circuit auto-recovers → attempts re-fire
# ==========================================================================
@settings(
    max_examples=10,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow],
)
@given(prompt=st.sampled_from(_PROMPTS))
def test_p4_ttl_recovery_refires_attempts(prompt):
    gw = _HealthyGW(_VALID_PNG_B64)
    _IMAGE_GEN_ATTEMPTS.clear()
    # Force-trip the circuit, then roll disabled_at back beyond the TTL window
    # (ttl=300s) so the next check sees it as expired.
    _IMAGE_GEN_CIRCUIT["disabled_at"] = time.time()
    _IMAGE_GEN_CIRCUIT["disabled_at"] = time.time() - 301

    # TTL elapsed → circuit reports recovered (and resets disabled_at to 0).
    assert _image_gen_is_circuit_broken() is False, (
        "circuit must auto-recover once TTL has elapsed"
    )

    saved = {
        k: os.environ.get(k)
        for k in ("AE_PREFER_VERTEX_IMAGE", "AE_IMAGE_QUALITY_THRESHOLD")
    }
    os.environ["AE_PREFER_VERTEX_IMAGE"] = "0"      # skip vertex-first path
    os.environ["AE_IMAGE_QUALITY_THRESHOLD"] = "0"  # no quality-retry recursion
    try:
        with tempfile.TemporaryDirectory() as tmp, \
                patch.object(srv, "_get_gw", return_value=gw), \
                patch.object(srv, "_resolve_callable_model_id",
                             side_effect=lambda m, *a, **k: m), \
                patch.object(srv, "_try_vertex_image_single", new=_async_none), \
                patch("ai_engine.vertex_image_module.get_vertex_image_client",
                      return_value=_DisabledVertexClient()):
            raw = asyncio.run(srv._tool_generate_image(
                {"prompt": prompt, "size": "1024x1024"},
                project_path=tmp,
            ))
        result = json.loads(raw)
        # After recovery the chain is retried — image models are invoked again.
        image_calls = [m for m in gw.calls if _is_image_model_id(m)]
        assert image_calls, (
            "after TTL recovery, _tool_generate_image must re-fire image-model "
            f"attempts; gw.calls={gw.calls!r}"
        )
        # And with a healthy gateway the result is a real image, not a circuit
        # short-circuit error.
        assert "error" not in result, (
            f"healthy gateway after recovery should succeed, got: {result!r}"
        )
    finally:
        for k, v in saved.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v
        _IMAGE_GEN_CIRCUIT["disabled_at"] = 0
        _IMAGE_GEN_ATTEMPTS.clear()


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
