"""산출물 audit 통합 게이트 테스트 — pptx-native-density-render (작업 18.1).

이 파일은 스펙의 **합격 게이트**(design.md §Testing Strategy "합격 게이트 정의 Req 8")
를 구현한다. hermetic 단위 통과가 아니라 **실제 생성된 .pptx 를 산출물_검증기로
audit** 하여 (a)~(h) 항목 전수 통과를 합격 조건으로 한다(Req 8.7).

검증 대상(실제 코드, 네트워크 0 · Chrome 미사용):
  - ``ai_engine.native_layout_renderer.render_native_layout`` 으로 표지(cover) +
    7개 알려진_레이아웃(cover/section_divider/two_column/feature_grid/timeline/
    comparison/architecture)을 **충분한 콘텐츠**로 채운 대표 덱을 실제 .pptx
    (13.333×7.5, 16:9)로 생성한다.
  - ``scripts.audit_pptx_native_density.audit_native_density`` 로 (a)~(h) 전수 통과
    ((h) 스타일 품질 ``audit_style_quality`` 포함)를 단언한다 — ``passed=True``
    이어야 하고, 미통과 시 ``failures`` 가 비어있지 않다(=불합격).
  - **Vertex 옵트인 ON(mock)/OFF 두 변형** 모두 생성·audit 한다. ON 은
    ``native_layout_renderer._get_vertex_client`` 를 monkeypatch 한 가짜 클라이언트로
    추상 PNG 를 반환(실제 Vertex/네트워크 호출 0)하고, 장식_배경을 풀블리드로
    임베드한 상태로 audit 가 (a)~(h) 를 통과하는지 확인한다(Req 11.5/11.6).
  - (i) ``generate_visual_comparison`` 산출물은 육안 보조이며 자동판정 비포함 —
    입력 누락 시 graceful skip(None) 임을 확인한다(Chrome 미가용/타임아웃/입력 누락
    시 예외 전파 없이 skip, 합격 판정에 영향 없음).

실행 규약(design.md §Testing Strategy):
  - **네트워크 0 / Chrome 미사용**: 게이트웨이/Vertex 실호출 없음(ON 은 mock).
    audit 는 python-pptx 만으로 동작(Chrome 불필요).
  - 임시 .pptx 는 OS 임시 경로에 쓰고 ``finally`` 에서 정리(잔존 0).
  - 인라인 ``-c``/heredoc 금지 — 파일로 작성, 단발 실행.

실행:
  ./venv/bin/python -m pytest scripts/test_pptx_native_density_audit_integration.py \
      -p no:cacheprovider -q
"""
from __future__ import annotations

import os
import io
import sys
import base64
import shutil
import tempfile
import contextlib
from unittest.mock import patch

# --- import 경로: repo 루트(ai_engine 패키지) + scripts 디렉터리(audit 모듈) ---
_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.dirname(_THIS_DIR)
_AE_DIR = os.path.join(_REPO_ROOT, "ai_engine")
for _p in (_REPO_ROOT, _AE_DIR, _THIS_DIR):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

# NOTE(모듈 단일성): Vertex mock 은 _get_vertex_client 를 monkeypatch 하므로,
# 덱을 *빌드하는* render_native_layout 과 patch 대상은 **동일 모듈 객체**여야 한다.
# (top-level ``native_layout_renderer`` 와 ``ai_engine.native_layout_renderer`` 는
#  sys.modules 상 별개 객체가 될 수 있어 patch 가 안 먹는 함정이 있다.) 따라서
# 단일 참조 ``nlr`` 만 쓰고 ``nlr.render_native_layout`` 를 호출한다.
import native_layout_renderer as nlr  # noqa: E402  (ai_engine on path)
from slide_templates import design_tokens_for_profile  # noqa: E402

import audit_pptx_native_density as aud  # noqa: E402  (scripts on path)
from audit_pptx_native_density import (  # noqa: E402
    audit_native_density,
    audit_style_quality,
    AuditReport,
    generate_visual_comparison,
    CHECK_EDITABLE_TEXT,
    CHECK_OVERLAP,
    CHECK_BOUNDS,
    CHECK_TITLE_ONCE,
    CHECK_DENSITY,
    CHECK_BAKED_TEXT,
    CHECK_ZORDER,
    CHECK_STYLE_QUALITY,
)

SLIDE_W = 13.333
SLIDE_H = 7.5

_TOKENS = design_tokens_for_profile(None)

# (a)~(h) 검사 항목 식별자 — 통합 게이트는 이 전수가 failures 에 없어야 합격.
_ALL_CHECKS = [
    CHECK_EDITABLE_TEXT, CHECK_OVERLAP, CHECK_BOUNDS, CHECK_TITLE_ONCE,
    CHECK_DENSITY, CHECK_BAKED_TEXT, CHECK_ZORDER, CHECK_STYLE_QUALITY,
]


# ===========================================================================
# 대표 덱 — 표지(cover) + 7개 알려진_레이아웃, 충분한 콘텐츠
# ===========================================================================

def _deck_specs():
    """표지 + 7개 알려진_레이아웃을 충분한 콘텐츠로 채운 대표 덱 명세."""
    return [
        ("cover", {
            "eyebrow": "2026 사업 전략",
            "title": "차세대 AI 플랫폼 통합 전략 보고서",
            "subtitle": "편집가능 네이티브 렌더링으로 구현한 젠스파크급 고밀도 디자인",
            "footer": "전략기획팀 · 2026-06 · Confidential",
        }),
        ("section_divider", {
            "section_number": 1,
            "title": "시장 환경 분석",
            "description": "글로벌 AI 시장의 구조적 변화와 우리의 대응 방향을 정리한다. "
                           "수요 측면의 폭발적 성장과 공급 측면의 경쟁 심화를 동시에 조망한다.",
        }),
        ("two_column", {
            "title": "현황과 과제",
            "subtitle": "내부 역량과 외부 환경의 균형 점검",
            "left_content": ["강력한 데이터 자산 보유", "검증된 추론 파이프라인",
                             "안정적 게이트웨이 인프라", "다년간 축적된 도메인 지식"],
            "right_content": ["멀티모달 대응 지연", "추론 비용 최적화 필요",
                              "거버넌스 체계 보강", "글로벌 규제 대응 강화"],
        }),
        ("feature_grid", {
            "title": "핵심 기능 4종",
            "subtitle": "플랫폼이 제공하는 차별화 역량",
            "features": [
                {"icon": "★", "title": "고밀도 렌더링", "description": "젠스파크급 밀도를 네이티브 도형으로 구현"},
                {"icon": "◆", "title": "편집 가능", "description": "모든 콘텐츠 텍스트가 PowerPoint 에서 편집 가능"},
                {"icon": "●", "title": "겹침 0", "description": "기하 보정으로 셰이프 겹침 10% 미만 보장"},
                {"icon": "▲", "title": "장식 분리", "description": "Vertex 장식은 콘텐츠 텍스트와 z-order 분리"},
            ],
        }),
        ("timeline", {
            "title": "구축 로드맵",
            "subtitle": "4단계 추진 일정",
            "steps": [
                {"label": "1분기", "title": "기반 구축", "description": "네이티브 렌더러 스캐폴드 완성"},
                {"label": "2분기", "title": "요소 매핑", "description": "고밀도 요소 도형화 및 토큰 적용"},
                {"label": "3분기", "title": "검증 게이트", "description": "산출물 audit 통합 게이트 구축"},
                {"label": "4분기", "title": "정식 출시", "description": "전사 배포 및 안정화 운영"},
            ],
        }),
        ("comparison", {
            "title": "기존 방식 대비 개선",
            "subtitle": "통짜 이미지 vs 네이티브 렌더",
            "left_label": "기존(통짜 이미지)",
            "left_items": ["편집 불가", "텍스트 픽셀 베이크", "겹침 잔존", "접근성 미흡"],
            "right_label": "신규(네이티브)",
            "right_items": ["완전 편집 가능", "텍스트 런 보존", "겹침 0 보정", "스크린리더 호환"],
        }),
        ("architecture", {
            "title": "시스템 아키텍처",
            "subtitle": "레이어 구성과 책임 분리",
            "layers": [
                {"name": "프레젠테이션 레이어", "description": "네이티브 도형 렌더", "items": ["emit_*", "design_tokens"]},
                {"name": "기하 보정 레이어", "description": "겹침/경계 보정", "items": ["resolve_collisions", "clamp"]},
                {"name": "검증 레이어", "description": "산출물 audit 게이트", "items": ["audit_native_density"]},
            ],
        }),
    ]


def _build_deck(prs, *, vertex_on: bool):
    """대표 덱을 render_native_layout 로 실제 슬라이드에 생성한다.

    vertex_on=True 면 각 슬라이드 data 에 추상 장식 힌트(figure_prompt)를 주입해
    Vertex 장식_배경 경로(mock)를 폭넓게 exercise 한다(콘텐츠 텍스트는 프롬프트에
    넣지 않음 → 베이크 금지).
    """
    for layout, data in _deck_specs():
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        d = dict(data)
        if vertex_on:
            # 추상 장식_배경 힌트(콘텐츠 텍스트 아님). cover/section_divider 는 기본
            # 프롬프트가 있으나, body 레이아웃도 장식 경로를 타도록 명시 주입한다.
            d["figure_prompt"] = ("abstract decorative gradient background, "
                                  "no text, no letters")
        nlr.render_native_layout(slide, prs, layout, d, _TOKENS)


def _new_prs():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


@contextlib.contextmanager
def _temp_pptx(prs):
    """Presentation 을 임시 .pptx 로 저장하고 경로를 yield, finally 에서 정리."""
    fd, path = tempfile.mkstemp(suffix=".generated.pptx", prefix="it_native_density_")
    os.close(fd)
    try:
        prs.save(path)
        yield path
    finally:
        try:
            os.remove(path)
        except OSError:
            pass


# ===========================================================================
# Vertex mock — 실제 Vertex/네트워크 호출 절대 금지 (추상 PNG)
# ===========================================================================

def _abstract_png(seed: int = 7) -> bytes:
    """장식_배경용 추상(평탄 단색) PNG — baked_text_score 판정 미초과(텍스트 없음)."""
    img = Image.new("RGB", (320, 200),
                    (40 + seed % 180, 60 + (seed * 3) % 160, 90 + (seed * 7) % 140))
    buf = io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


class _FakeVertexClient:
    """Vertex 이미지 클라이언트 mock — 네트워크 0. async generate 로 추상 PNG 반환."""

    def __init__(self, png_bytes: bytes, *, enabled: bool = True):
        self._png = png_bytes
        self.enabled = enabled
        self.calls = 0

    async def generate(self, prompt=None, model_class=None, aspect_ratio=None, **_kw):
        self.calls += 1
        return {"images": [base64.b64encode(self._png).decode("ascii")]}


@contextlib.contextmanager
def _optin_env(value):
    key = "AE_ENABLE_VERTEX_IMAGE"
    old = os.environ.get(key)
    if value is None:
        os.environ.pop(key, None)
    else:
        os.environ[key] = value
    try:
        yield
    finally:
        if old is None:
            os.environ.pop(key, None)
        else:
            os.environ[key] = old


@contextlib.contextmanager
def _vertex_on(fake):
    """_get_vertex_client 를 fake 로 patch + 옵트인 ON(컨텍스트). 실제 호출 0."""
    with patch.object(nlr, "_get_vertex_client", lambda **_kw: fake), _optin_env("1"):
        yield


def _clean_decor_tmp():
    shutil.rmtree(os.path.join(tempfile.gettempdir(), "ae_vertex_decor"),
                  ignore_errors=True)


# ===========================================================================
# 슬라이드 트리 헬퍼
# ===========================================================================

def _count_pictures(prs) -> int:
    n = 0
    for slide in prs.slides:
        for sh in slide.shapes:
            try:
                if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    n += 1
            except Exception:
                pass
    return n


def _fail_summary(report: AuditReport) -> str:
    return "; ".join(
        f"[slide {f['slide']}] {f['check']}: {f['signal']}" for f in report.failures
    ) or "(없음)"


# ===========================================================================
# 합격 게이트 1 — Vertex 옵트인 OFF 변형: (a)~(h) 전수 통과
# ===========================================================================

def test_audit_gate_vertex_off_passes_all_checks():
    """대표 덱(표지+7개 레이아웃)을 Vertex OFF 로 실제 .pptx 생성 → audit_native_density
    가 (a)~(h) 전수 통과(passed=True, failures 없음)함을 단언한다.

    합격은 hermetic 단위가 아니라 **실제 산출물 audit** 으로 정의된다(Req 8.7).

    **Validates: Requirements 8.1, 8.2, 8.3, 8.4, 8.5, 8.6, 8.8, 8.9, 8.10, 10.1, 10.2**
    """
    with _optin_env("0"):  # Vertex OFF(네트워크 0, 장식 이미지 미생성)
        prs = _new_prs()
        _build_deck(prs, vertex_on=False)
        # OFF 변형: 콘텐츠 텍스트가 통짜 이미지로 베이크되지 않음(그림 0).
        assert _count_pictures(prs) == 0, "Vertex OFF 인데 그림이 임베드됨(베이크 의심)"

        with _temp_pptx(prs) as path:
            report = audit_native_density(path, _TOKENS)

    # (a)~(h) 전수 통과 = 합격.
    assert isinstance(report, AuditReport)
    assert report.passed is True, f"audit 미통과(불합격): {_fail_summary(report)}"
    assert report.failures == [], f"failures 존재: {_fail_summary(report)}"

    # 7개 슬라이드 전부 검사됐고, 어떤 (a)~(h) 항목도 실패하지 않음.
    failed_checks = {f["check"] for f in report.failures}
    for chk in _ALL_CHECKS:
        assert chk not in failed_checks, f"{chk} 실패"


# ===========================================================================
# 합격 게이트 2 — Vertex 옵트인 ON(mock) 변형: (a)~(h) 전수 통과 + 장식 임베드
# ===========================================================================

def test_audit_gate_vertex_on_mock_passes_all_checks():
    """대표 덱을 Vertex ON(mock) 으로 실제 .pptx 생성 → 장식_배경(추상 PNG)을
    풀블리드로 임베드한 상태에서도 audit_native_density 가 (a)~(h) 전수 통과함을
    단언한다. 실제 Vertex/네트워크 호출은 0(_get_vertex_client monkeypatch).

    **Validates: Requirements 8.1, 8.2, 8.3, 8.4, 8.5, 8.6, 8.8, 8.9, 8.10, 11.5, 11.6**
    """
    fake = _FakeVertexClient(_abstract_png(11), enabled=True)
    _clean_decor_tmp()
    try:
        with _vertex_on(fake):
            prs = _new_prs()
            _build_deck(prs, vertex_on=True)

            # ON 변형: 장식_배경(추상 PNG)이 실제로 임베드됨(풀블리드 장식 슬롯).
            n_pics = _count_pictures(prs)
            assert n_pics >= 1, "Vertex ON 인데 장식 이미지가 임베드되지 않음"
            # 실제 Vertex 호출 0 — mock(_get_vertex_client) 만 사용됨(네트워크 0).
            assert fake.calls >= 1, "mock Vertex 클라이언트가 호출되지 않음(경로 미실행)"

            with _temp_pptx(prs) as path:
                report = audit_native_density(path, _TOKENS)
    finally:
        _clean_decor_tmp()

    # (f) 베이크 텍스트: 추상 장식_배경은 판정 미초과여야 함(콘텐츠 공존, Req 11.6).
    assert report.passed is True, f"ON 변형 audit 미통과: {_fail_summary(report)}"
    assert report.failures == [], f"ON 변형 failures: {_fail_summary(report)}"
    failed_checks = {f["check"] for f in report.failures}
    assert CHECK_BAKED_TEXT not in failed_checks, "장식_배경이 베이크 텍스트로 오검출"
    assert CHECK_ZORDER not in failed_checks, "콘텐츠 텍스트가 장식 이미지 아래 z-order"


# ===========================================================================
# 합격 게이트 3 — 슬라이드별 (a)~(h) 직접 확인 (보고 완전성)
# ===========================================================================

def test_audit_gate_per_slide_all_layouts_pass():
    """대표 덱의 표지+7개 레이아웃 각각이 audit (a)~(h) 와 audit_style_quality(h)를
    개별적으로 통과함을 슬라이드 단위로 단언한다(Req 8 합격 게이트, Property 18 공존).

    **Validates: Requirements 8.5, 5.1, 5.3, 5.4, 11.6**
    """
    specs = _deck_specs()
    with _optin_env("0"):
        prs = _new_prs()
        _build_deck(prs, vertex_on=False)
        with _temp_pptx(prs) as path:
            report = audit_native_density(path, _TOKENS)
            # 슬라이드별 스타일 품질(h) 개별 통과 확인.
            prs2 = Presentation(path)
            for i, slide in enumerate(prs2.slides):
                rep = audit_style_quality(slide, _TOKENS)
                layout = specs[i][0]
                assert rep.passed, (f"[slide {i+1} {layout}] 스타일 품질 미달: "
                                    f"{rep.missing_style} (score={rep.score:.2f})")

    assert report.passed is True, f"통합 audit 미통과: {_fail_summary(report)}"
    # 표지(cover)는 7개 알려진_레이아웃에 포함되므로 대표 덱은 7개 슬라이드다.
    assert len(specs) == 7


# ===========================================================================
# Req 8.7 — hermetic 단위 통과만으로는 합격 불가 (audit 미통과 시 불합격)
# ===========================================================================

def test_hermetic_unit_pass_alone_is_insufficient():
    """합격 게이트는 실제 .pptx 산출물 audit 이며, hermetic 단위 수준의 빈약한
    슬라이드(밀도/스타일 미달)는 audit 에서 **불합격**(passed=False)함을 단언한다.

    즉 "audit 미통과 시 실패" 구조로 Req 8.7(단위 통과만으로 합격 불가)을 표현한다.

    **Validates: Requirements 8.6, 8.7**
    """
    with _optin_env("0"):
        prs = _new_prs()
        # render_native_layout 을 거치지 않은 빈약한 슬라이드(평문 텍스트박스 1개) —
        # 단위 수준으로는 "텍스트가 있다" 정도만 충족하나 밀도/스타일 게이트는 미달.
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(6.0), Inches(1.0))
        tb.text_frame.text = "단위 수준 텍스트만 있는 빈약한 슬라이드"

        with _temp_pptx(prs) as path:
            report = audit_native_density(path, _TOKENS)

    # 빈약한 슬라이드는 합격 게이트를 통과하지 못한다(밀도(e)/스타일(h) 등 미달).
    assert report.passed is False, "빈약한 단위 슬라이드가 audit 합격(게이트 무력)"
    assert len(report.failures) >= 1
    failed = {f["check"] for f in report.failures}
    # 최소한 밀도(e) 또는 스타일 품질(h) 항목이 불합격이어야 한다.
    assert (CHECK_DENSITY in failed) or (CHECK_STYLE_QUALITY in failed), (
        f"밀도/스타일 게이트가 작동하지 않음: {failed}")


# ===========================================================================
# (i) visual_comparator 산출물 — 육안 보조, 자동판정 비포함(skip 가능)
# ===========================================================================

def test_visual_comparison_is_optional_and_skips_gracefully():
    """(i) generate_visual_comparison 은 자동판정((a)~(h))에 포함되지 않는 육안
    보조 산출물이며, 입력 누락/Chrome 미가용/타임아웃 시 예외 없이 skip(None) 한다.

    합격 자동판정은 audit_native_density(python-pptx 만)로만 수행되므로, 본 함수의
    skip 여부는 합격 결과에 영향을 주지 않는다(Req 5.1).

    **Validates: Requirements 5.1**
    """
    # 입력 누락(빈 HTML) → graceful skip(None), 예외 전파 없음.
    assert generate_visual_comparison("", "/nonexistent/ref.png", timeout=5.0) is None
    # 존재하지 않는 참조 PNG → skip(None).
    assert generate_visual_comparison("<html><body>x</body></html>",
                                      "/nonexistent/ref.png", timeout=5.0) is None
