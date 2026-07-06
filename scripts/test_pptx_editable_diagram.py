"""Regression — 편집 가능 네이티브 다이어그램이 Vertex 이미지보다 우선된다(젠스파크 방식).

사용자 핵심 요구: "젠스파크는 PPT의 모든 도형이 PowerPoint에서 편집 가능하다."
통짜 래스터 이미지(Vertex/HTML)는 편집 불가 → 거부. 따라서 섹션을 LLM으로 구조화해
네이티브 python-pptx 도형(편집 가능)으로 그려야 한다.

수정: _tool_generate_pptx가 Vertex 일괄 생성 *이전에* _llm_structure_native_diagram으로
각 섹션을 {type, content} 스펙으로 구조화해 sd["nativeDiagram"]에 기록 →
build_native_diagram이 편집 가능 도형(ROUNDED_RECTANGLE + connector)으로 렌더.

Correctness properties:
  P1. 게이트웨이가 flow 스펙을 주면 슬라이드에 편집 가능 도형(AutoShape)이 그려진다.
  P2. 네이티브 다이어그램이 그려진 슬라이드에는 Vertex 래스터 이미지(Picture)가 임베드되지 않는다.
  P3. 구조화 실패(빈 스펙) 시 크래시 없이 생성 완료(다른 경로로 폴백).

실행: pytest scripts/test_pptx_editable_diagram.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
vmod = pytest.importorskip("ai_engine.vertex_image_module")
pptx = pytest.importorskip("pptx")


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def _enable_struct(monkeypatch, spec):
    """게이트웨이/모델/구조화 함수를 mock해 항상 주어진 spec을 반환."""
    monkeypatch.setattr(server, "_get_gw", lambda *a, **k: object())
    monkeypatch.setattr(server, "_specialized_model_for_task", lambda *a, **k: "model-x")
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)

    async def _fake_struct(gw, model, h, b, ctx=""):
        return dict(spec) if spec else {}
    monkeypatch.setattr(server, "_llm_structure_native_diagram", _fake_struct)


def test_editable_native_diagram_drawn(tmp_path, monkeypatch):
    """P1+P2 — flow 스펙이면 편집 가능 도형이 그려지고 Vertex 이미지는 없다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "1")
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "1")
    # 이 테스트는 네이티브 편집 가능 경로 검증 — HTML 고품질 렌더는 끈다.
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "0")
    _enable_struct(monkeypatch, {"type": "flow", "content": "수집 -> 정제 -> 적재 -> 시각화"})

    # Vertex가 호출되면 실패시켜 두되, 애초에 nativeDiagram이 set되면 skip되어야 한다.
    class _Vtx:
        enabled = True
        async def generate(self, **k):
            raise AssertionError("Vertex must not run when native diagram is structured")
    monkeypatch.setattr(vmod, "get_vertex_image_client", lambda *a, **k: _Vtx())

    slides = [{"title": "데이터 처리 흐름", "bullets": ["원천에서 수집", "정제 후 적재", "대시보드 시각화"]}]
    out = asyncio.run(server._tool_generate_pptx({"title": "프로젝트 흐름도", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        content_slide = list(prs.slides)[1]
        # P1 — 편집 가능 AutoShape(둥근 사각형) 다수 존재
        autoshapes = [sh for sh in content_slide.shapes
                      if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        assert len(autoshapes) >= 3, f"편집 가능 도형 수={len(autoshapes)} (flow 4노드 기대)"
        # P2 — Vertex 래스터 이미지(Picture) 없음
        pics = [sh for sh in content_slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert not pics, "네이티브 다이어그램 슬라이드에 래스터 이미지가 임베드됨(편집 불가)"
    finally:
        _cleanup(res)


def test_struct_failure_safe_fallback(tmp_path, monkeypatch):
    """P3 — 구조화 실패(빈 스펙) 시 크래시 없이 생성 완료."""
    from pptx import Presentation
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "1")
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "0")  # Vertex 끔 → 네이티브 분류기/텍스트 폴백
    # 네이티브 경로 검증 — HTML 고품질 렌더는 끈다.
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "0")
    _enable_struct(monkeypatch, {})  # 빈 스펙

    slides = [{"title": "개요", "bullets": ["항목 1", "항목 2"]}]
    out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        assert len(prs.slides._sldIdLst) == 2
    finally:
        _cleanup(res)


def test_diagram_and_vertex_coexist(tmp_path, monkeypatch):
    """편집 가능 다이어그램 + Vertex 이미지 공존(역할 분담, 젠스파크 방식).

    구조화가 다이어그램 섹션(흐름)을 편집 가능 도형으로 claim하고,
    비다이어그램 섹션(개요)은 Vertex Nano Banana Pro 이미지로 채워진다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "1")
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "1")
    # 이 테스트는 네이티브 다이어그램 경로 검증 — HTML 고품질 렌더는 끈다.
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "0")
    # 공유 배경(hero/body bg)은 별도 기능 — 이 테스트는 다이어그램 vs 슬라이드별
    # Vertex 이미지 역할분담을 검증하므로 배경은 끈다(흐름 슬라이드 래스터 0 유지).
    monkeypatch.setenv("AE_DISABLE_VERTEX_HERO", "1")
    monkeypatch.setenv("AE_DISABLE_VERTEX_BODY_BG", "1")
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)
    monkeypatch.setattr(server, "_get_gw", lambda *a, **k: object())
    monkeypatch.setattr(server, "_specialized_model_for_task", lambda *a, **k: "model-x")

    # 구조화: '흐름' 섹션만 flow 스펙, 나머지는 비다이어그램({})
    async def _struct(gw, model, h, b, ctx=""):
        if "흐름" in (h or ""):
            return {"type": "flow", "content": "수집 -> 정제 -> 적재"}
        return {}
    monkeypatch.setattr(server, "_llm_structure_native_diagram", _struct)

    # Vertex: 활성 + PNG 생성(비다이어그램 섹션용)
    _PNG = bytes.fromhex(
        "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
        "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082")
    import base64 as _b64
    _b64png = _b64.b64encode(_PNG).decode()

    class _Vtx:
        enabled = True
        async def generate(self, **k):
            return {"images": [_b64png], "model": "gemini-3-pro-image-preview"}
    monkeypatch.setattr(vmod, "get_vertex_image_client", lambda *a, **k: _Vtx())

    slides = [
        {"title": "데이터 처리 흐름", "bullets": ["수집", "정제", "적재"]},
        {"title": "프로젝트 개요", "bullets": ["배경과 목표 설명"]},
    ]
    out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        sl_flow, sl_overview = list(prs.slides)[1], list(prs.slides)[2]
        # 흐름 슬라이드: 편집 가능 도형 + 래스터 없음
        flow_auto = [s for s in sl_flow.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        flow_pic = [s for s in sl_flow.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert len(flow_auto) >= 3 and not flow_pic, "흐름 슬라이드는 편집 가능 도형이어야 함"
        # 개요 슬라이드: Vertex 이미지(Picture) 존재
        ov_pic = [s for s in sl_overview.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        assert ov_pic, "비다이어그램 슬라이드에 Vertex 이미지가 들어가야 함"
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
