"""Regression — 모델이 generate_pptx를 직접 호출해도 다이어그램형 슬라이드는
네이티브(편집 가능) 도형으로 그려진다 (통짜 PNG 금지).

대상: ai_engine.server._tool_generate_pptx
배경: 이전엔 nativeDiagram은 강제 생성 폴백에서만 주입돼, 모델 직접 호출 경로는
      imagePrompt→통짜 PNG였다. 이제 직접 호출도 다이어그램형이면 네이티브 도형.

검증:
  - 다이어그램형 슬라이드(흐름/구조/아키텍처 bullets) → autoshape(도형) ≥ 1,
    embedded picture = 0 (편집 가능)
  - 의미성: AE_DISABLE_NATIVE_DIAGRAM=1이면 네이티브 비활성(폴백)

실행: pytest scripts/test_pptx_direct_native_diagram.py -q
"""
from __future__ import annotations

import asyncio
import os
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402


def _gen(tmp_root, slides):
    os.environ["AE_GENERATED_ROOT"] = tmp_root
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "흐름 테스트", "slides": slides}, tmp_root))
    # JSON 결과의 absPath를 신뢰(상대경로 재구성은 _resolve_local_root 정책에 취약).
    import json as _json
    try:
        data = _json.loads(out)
    except (ValueError, TypeError):
        data = {}
    path = data.get("absPath") or data.get("path") or ""
    assert path, f"pptx 경로를 찾지 못함: {out[:300]}"
    if not os.path.isabs(path):
        path = os.path.join(tmp_root, path)
    assert os.path.isfile(path), f"생성 파일 없음: {path}"
    return path


def _count_shapes(path):
    prs = Presentation(path)
    autoshapes = 0
    pictures = 0
    connectors = 0
    for slide in prs.slides:
        for shp in slide.shapes:
            st = shp.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                autoshapes += 1
            elif st == MSO_SHAPE_TYPE.LINE:
                connectors += 1
    return autoshapes, pictures, connectors


# 다이어그램형 — imagePrompt 없이 구조/흐름 bullets만. 게이트웨이 호출 불필요.
_DIAGRAM_SLIDES = [{
    "title": "시스템 아키텍처",
    "bullets": ["프레젠테이션 계층", "애플리케이션 계층", "데이터 계층"],
}]


def test_direct_call_produces_native_shapes():
    with tempfile.TemporaryDirectory() as td:
        os.environ.pop("AE_DISABLE_NATIVE_DIAGRAM", None)
        path = _gen(td, [dict(s) for s in _DIAGRAM_SLIDES])
        autoshapes, pictures, _conn = _count_shapes(path)
        assert autoshapes >= 1, "네이티브 도형이 없음 — 통짜 이미지로 생성됨"
        assert pictures == 0, "다이어그램형 슬라이드에 embedded picture가 들어감(편집 불가)"


def test_disable_flag_falls_back():
    # 의미성 증명 — 비활성 시 네이티브 도형을 강제로 그리지 않는다.
    with tempfile.TemporaryDirectory() as td:
        os.environ["AE_DISABLE_NATIVE_DIAGRAM"] = "1"
        try:
            path = _gen(td, [dict(s) for s in _DIAGRAM_SLIDES])
            autoshapes, _pic, _conn = _count_shapes(path)
            # 비활성 시 분류기 추론 경로를 타지 않으므로 다이어그램 도형이 생기지 않는다
            # (본문 placeholder 텍스트만). autoshape가 0이거나 매우 적어야 한다.
            assert autoshapes == 0
        finally:
            os.environ.pop("AE_DISABLE_NATIVE_DIAGRAM", None)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
