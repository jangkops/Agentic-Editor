"""아이콘 + 설명 디테일 회귀.

- 아이콘 키워드 매핑(의미 키) 검증
- 리치 아이콘 비활성(AE_DISABLE_RICH_ICONS=1): 카드 PICTURE=0, autoshape 폴백
- 리치 아이콘 활성(Chrome 존재 시): Lucide 벡터 아이콘 PNG 임베드
- block 제목+설명, flow 설명 유지 검증
"""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "ai_engine"))
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import native_diagram_pptx
from native_diagram_pptx import build_native_diagram, _icon_name_for
import icon_assets


def _stats(slide):
    pic = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    ash = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE)
    texts = " ".join(s.text_frame.text for s in slide.shapes
                     if s.has_text_frame and s.text_frame.text.strip())
    return pic, ash, texts


def test_icon_mapping():
    assert _icon_name_for("AWS 클라우드 서버") == "cloud"
    assert _icon_name_for("자동화 프로비저닝 파이프라인") == "gear"
    assert _icon_name_for("데이터 저장소 백업") == "database"
    assert _icon_name_for("보안 감사 권한 통제") == "shield"
    assert _icon_name_for("월 비용 예산 절감") == "dollar"
    assert _icon_name_for("실시간 성능 가속") == "zap"
    assert _icon_name_for("사용자 팀 계정 관리") in ("users", "gear")
    assert _icon_name_for("그냥 일반 항목") == "chart"
    print("[icon mapping] PASS")


CARDS = ("실시간 성능: 빠른 응답속도\n클라우드 인프라: AWS 확장\n보안 감사: 권한 통제\n"
         "비용 절감: 예산 최적화\n데이터 저장소: 자동 백업\n핵심 기능: 대표 가치")


def test_fallback_mode():
    os.environ["AE_DISABLE_RICH_ICONS"] = "1"
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    ok = build_native_diagram(s, "cards", CARDS, palette=["#4472C4", "#ED7D31", "#A5A5A5"])
    pic, ash, t = _stats(s)
    print(f"[fallback] ok={ok} pic={pic} autoshape={ash}")
    assert ok and pic == 0 and ash >= 18, f"fallback pic={pic} a={ash}"
    assert "빠른 응답속도" in t
    del os.environ["AE_DISABLE_RICH_ICONS"]


def test_rich_mode():
    if not icon_assets.available():
        print("[rich] SKIP (Chrome 없음 — 폴백만 사용)")
        return
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    ok = build_native_diagram(s, "cards", CARDS, palette=["#4472C4", "#ED7D31", "#A5A5A5"])
    pic, ash, t = _stats(s)
    print(f"[rich] ok={ok} pic(아이콘)={pic} autoshape={ash}")
    assert ok and pic >= 4, f"리치 아이콘 임베드 안 됨 pic={pic}"
    assert ash >= 12, "카드/칩 등 네이티브 도형 유지 확인"
    assert "빠른 응답속도" in t


def test_block_flow():
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    ok2 = build_native_diagram(s2, "block",
        "포털 사용자 관리: 리스트 확인\n통합 프로비저닝: 자동 설정\n서버 계정 생성: 원클릭",
        palette=["#4472C4", "#ED7D31"])
    _, _, t2 = _stats(s2)
    assert ok2 and "리스트 확인" in t2, "block 설명 누락"
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    ok3 = build_native_diagram(s3, "flow",
        "분석: 요구사항 -> 설계: 아키텍처 -> 개발: 구현 -> 배포: 운영", palette=["#4472C4"])
    _, _, t3 = _stats(s3)
    assert ok3 and "요구사항" in t3
    print("[block/flow] PASS")


def main():
    test_icon_mapping()
    test_fallback_mode()
    test_rich_mode()
    test_block_flow()
    print("ALL PASS")


if __name__ == "__main__":
    main()
