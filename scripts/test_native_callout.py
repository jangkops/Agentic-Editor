#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Regression — build_native_diagram의 note 파라미터가 NOTICE 콜아웃(크림 박스 +
주황 좌측 보더 + 'NOTICE' 라벨 + 본문)을 편집 가능한 네이티브 도형으로 그린다.

note 파라미터 부재 시 서버의 _bnd(note=...) 호출이 TypeError로 조용히 실패해
네이티브 경로가 PNG로 폴백하던 잠재 버그의 회귀 가드도 겸한다.

실행: python scripts/test_native_callout.py
"""
import os
import sys

_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
sys.path.insert(0, _ROOT)
sys.path.insert(0, os.path.join(_ROOT, "ai_engine"))

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from native_diagram_pptx import build_native_diagram


def _fill_hex(sh):
    try:
        return str(sh.fill.fore_color.rgb)
    except Exception:
        return ""


def test_callout():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    ok = build_native_diagram(
        s, "cards",
        "성능 최적화: 실시간 응답\n보안 강화: 권한 통제\n데이터 백업: 자동화",
        palette=["#4472C4", "#ED7D31", "#A5A5A5"],
        note="미설치 시 네트워크 접근이 차단됩니다. Teams로 담당자에게 문의하세요.",
    )
    assert ok, "build_native_diagram returned False"
    cream = sum(1 for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and _fill_hex(sh) == "FFF8EE")
    orange = sum(1 for sh in s.shapes
                 if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and _fill_hex(sh) == "ED7D31")
    notice = sum(1 for sh in s.shapes
                 if sh.has_text_frame and "NOTICE" in sh.text_frame.text)
    assert cream >= 1, f"cream callout box 없음 (cream={cream})"
    assert orange >= 1, f"orange 좌측 보더 없음 (orange={orange})"
    assert notice >= 1, f"NOTICE 라벨 없음 (notice={notice})"
    print(f"[callout] ok cream={cream} orange={orange} notice={notice}")


def test_no_note_no_callout():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    ok = build_native_diagram(
        s, "flow", "수집 -> 정제 -> 적재",
        palette=["#4472C4", "#ED7D31"])
    assert ok, "flow without note failed"
    cream = sum(1 for sh in s.shapes
                if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and _fill_hex(sh) == "FFF8EE")
    assert cream == 0, f"note 없는데 콜아웃 생성됨 (cream={cream})"
    print("[no-note] ok no callout")


if __name__ == "__main__":
    test_callout()
    test_no_note_no_callout()
    print("ALL PASS")
