"""Hermetic integration test — B방향(초고퀄 비주얼 우선, 통짜 이미지 하나).

목표(사용자 확정 B): **각 슬라이드 = 통짜 이미지 하나 + 그 위에 겹치는 네이티브
오버레이 0(겹침·중복 완전 제거)**. 편집가능 네이티브 라우팅(AE_NATIVE_LAYOUT_RENDER)
은 B와 상충하므로 기본 OFF다(미설정=OFF).

이 테스트는 실제 ``ai_engine/server.py:_tool_generate_pptx`` 결정 경로를 헤르메틱하게
구동한다(네트워크 없음). Chrome 감지·게이트웨이·HTML→PNG 렌더러·섹션 렌더러·이미지
생성은 모두 mock 이며, HTML 렌더는 **성공 mock**(프로덕션처럼 통짜 배경 PNG 경로를
반환하고 로컬에 실제 PNG 바이트를 기록)한다. Vertex 는 OFF(AE_PREFER_VERTEX_IMAGE=0).

혼합 덱(표지 + 본문 + 다이어그램[nativeDiagram 명시])을 생성한 뒤, 생성된 .pptx 를
직접 검사해 B 전용 기준을 단언한다(audit_native_density 재사용 안 함 — 그건 "편집가능
텍스트≥1"을 요구해 B와 상충):

  * 모든 슬라이드에서: 통짜 배경 PICTURE(0,0 에서 13.333×7.5in)가 있으면 그 위에
    겹치는(area overlap ≥10%) 네이티브 텍스트/도형(TEXT_BOX/AUTO_SHAPE, 또는 텍스트가
    있는 PLACEHOLDER) 셰이프 수 == 0.
  * 표지(index 0)도 위 조건 충족(수정1 검증 — coverBackground 통짜).
  * 다이어그램 슬라이드도 위 조건 충족(수정2 검증 — nativeDiagram + slideBackground
    → 통짜 배경만, 네이티브 다이어그램 억제).
  * 모든 도형이 슬라이드 경계(0,0,13.333,7.5)in 안(수정5).

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_pptx_bmode_solid_integration.py -p no:cacheprovider -q
"""
from __future__ import annotations

import io
import os
import sys
import json
import asyncio
import tempfile
from unittest.mock import patch

sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402


# --------------------------------------------------------------------------
# Hermetic fakes
# --------------------------------------------------------------------------
def _write_png(path: str, tag: int) -> None:
    os.makedirs(os.path.dirname(path), exist_ok=True)
    img = Image.new("RGB", (64 + (tag % 9), 48 + (tag % 5)),
                    (tag % 256, (tag * 7) % 256, (tag * 31) % 256))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    with open(path, "wb") as f:
        f.write(buf.getvalue())


async def _img_gen_disabled(*_a, **_k):
    return json.dumps({"error": "test-disabled"})


async def _render_html_png_ok(html, output_path, width=1920, height=1080, timeout=30, **_k):
    """_render_html_slide_to_png 성공 mock — 실제 PNG 바이트를 기록."""
    _write_png(output_path, 900_000)
    return {"ok": True}


def _make_section_bg_fake(project_path: str):
    """_generate_html_slide_for_section mock — 본문 슬라이드용 통짜 배경 rel 반환.

    프로덕션처럼 통짜 배경 PNG 경로를 반환한다 → 호출부가 sd['slideBackground'] 로
    채택해 풀블리드로 임베드한다(콘텐츠가 구워진 통짜)."""
    state = {"n": 0}

    async def _section_fake(gw, model_id, heading, body, doc_context,
                            proj_path, style_profile=None,
                            hero_image="", render_info=None):
        state["n"] += 1
        name = f"section-bg-{state['n']}.png"
        out = os.path.join(project_path, ".generated", name)
        _write_png(out, 700_000 + state["n"])
        if isinstance(render_info, dict):
            render_info["layout"] = "feature_grid"
            render_info["composited"] = False
        return f".generated/{name}"

    return _section_fake


# --------------------------------------------------------------------------
# PPTX inspection helpers (B 전용 — audit_native_density 재사용 안 함)
# --------------------------------------------------------------------------
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_TOL = Inches(0.05)


def _bbox(sh):
    return (sh.left or 0, sh.top or 0, sh.width or 0, sh.height or 0)


def _fullbleed_pictures(slide):
    pics = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        l, t, w, h = _bbox(sh)
        if (abs(l - 0) <= _TOL and abs(t - 0) <= _TOL
                and abs(w - _SLIDE_W) <= _TOL and abs(h - _SLIDE_H) <= _TOL):
            pics.append(sh)
    return pics


def _overlap_ratio(sh, pic):
    """intersection_area / shape_area (shape 자기 면적 대비 겹침 비율)."""
    sl, st, sw, sh_ = _bbox(sh)
    pl, pt, pw, ph = _bbox(pic)
    ix = max(0, min(sl + sw, pl + pw) - max(sl, pl))
    iy = max(0, min(st + sh_, pt + ph) - max(st, pt))
    inter = ix * iy
    area = sw * sh_
    return (inter / area) if area > 0 else 0.0


def _shape_text(sh) -> str:
    try:
        if sh.has_text_frame:
            return (sh.text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def _overlapping_native_overlays(slide):
    """통짜 배경 위에 겹치는(≥10%) 네이티브 텍스트/도형 offenders 목록.

    offender = 비-PICTURE shape 이면서 배경과 면적 겹침 ≥10% 이고,
      (shape_type in {TEXT_BOX, AUTO_SHAPE} OR 텍스트가 있는 PLACEHOLDER/기타).
    통짜 배경 PICTURE 자체는 제외. 배경이 없으면 [](공존 없음)."""
    pics = _fullbleed_pictures(slide)
    if not pics:
        return []
    offenders = []
    for sh in slide.shapes:
        try:
            st = sh.shape_type
        except Exception:
            st = None
        if st == MSO_SHAPE_TYPE.PICTURE:
            continue
        txt = _shape_text(sh)
        is_native_visual = st in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE)
        is_text_bearing = bool(txt)
        if not (is_native_visual or is_text_bearing):
            continue
        if any(_overlap_ratio(sh, pic) >= 0.10 for pic in pics):
            offenders.append((str(st), txt[:40]))
    return offenders


def _out_of_bounds(slide):
    """경계(0,0,13.333,7.5)in 밖 shape 목록(작은 tolerance 허용)."""
    bad = []
    tol = Inches(0.02)
    for sh in slide.shapes:
        l, t, w, h = _bbox(sh)
        if l < -tol or t < -tol or (l + w) > (_SLIDE_W + tol) or (t + h) > (_SLIDE_H + tol):
            bad.append((str(getattr(sh, "shape_type", "?")), l, t, w, h))
    return bad


# --------------------------------------------------------------------------
# 혼합 덱 생성 (표지 + 본문 + 다이어그램)
# --------------------------------------------------------------------------
def _generate_bmode_deck():
    proj = tempfile.mkdtemp()
    # 표지 통짜 배경(수정1) + 다이어그램 통짜 배경(수정2)용 실제 PNG 준비.
    _write_png(os.path.join(proj, ".generated", "cover-bg.png"), 111)
    _write_png(os.path.join(proj, ".generated", "diagram-bg.png"), 222)

    deck = {
        "title": "B방향 통짜 검증 덱",
        "coverBackground": ".generated/cover-bg.png",
        "slides": [
            {
                "title": "핵심 기능 요약",
                "bullets": ["빠른 처리", "안정적 운영", "유연한 확장",
                            "강력한 보안", "비용 절감", "쉬운 사용성"],
            },
            {
                "title": "업무 처리 프로세스",
                "nativeDiagram": {"type": "flow",
                                  "content": "접수 -> 검토 -> 승인 -> 완료"},
                "slideBackground": ".generated/diagram-bg.png",
            },
        ],
    }
    env = {
        # B방향: 네이티브 라우팅 기본 OFF — 명시적으로 미설정 상태를 보장(pop).
        "AE_ENABLE_HTML_SLIDES": "1",     # HTML 통짜 주 렌더러 ON
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",    # Vertex OFF (헤르메틱)
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    with patch.dict(os.environ, env, clear=False):
        os.environ.pop("AE_NATIVE_LAYOUT_RENDER", None)  # 기본 OFF 보장
        with patch.object(srv, "_call_bridge", lambda *a, **k: None), \
                patch.object(srv, "_find_local_chrome", lambda: "/fake/chrome"), \
                patch.object(srv, "_get_gw", lambda *a, **k: object()), \
                patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
                patch.object(srv, "_render_html_slide_to_png", _render_html_png_ok), \
                patch.object(srv, "_generate_html_slide_for_section",
                             _make_section_bg_fake(proj)), \
                patch.object(srv, "_tool_generate_image", _img_gen_disabled):
            raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"
    return result["absPath"]


# ==========================================================================
# 테스트
# ==========================================================================
def test_bmode_no_native_overlay_over_fullbleed_all_slides():
    """B 핵심: 모든 슬라이드에서 통짜 배경 위 겹치는 네이티브 오버레이 == 0."""
    pptx = _generate_bmode_deck()
    slides = list(Presentation(pptx).slides)
    assert len(slides) == 3, f"슬라이드 수 기대 3(표지+본문+다이어그램) — 실제 {len(slides)}"

    for idx, slide in enumerate(slides):
        offenders = _overlapping_native_overlays(slide)
        assert offenders == [], (
            f"슬라이드 {idx}: 통짜 배경 위 겹치는 네이티브 오버레이 {len(offenders)}개 "
            f"(겹침0·중복0 위반) — {offenders}")


def test_bmode_cover_is_solid_only():
    """수정1 검증: 표지(index 0)는 통짜 배경 1장 + 겹치는 네이티브 텍스트/도형 0."""
    pptx = _generate_bmode_deck()
    cover = list(Presentation(pptx).slides)[0]
    pics = _fullbleed_pictures(cover)
    assert len(pics) >= 1, "표지에 통짜(풀블리드) 배경 PICTURE 가 있어야 함(coverBackground 채택)"
    offenders = _overlapping_native_overlays(cover)
    assert offenders == [], (
        f"표지 통짜 배경 위 겹치는 네이티브 셰이프 {len(offenders)}개 (수정1 위반) — {offenders}")


def test_bmode_diagram_slide_is_solid_only():
    """수정2 검증: 다이어그램 슬라이드는 통짜 배경만 + 네이티브 다이어그램 억제(겹침0)."""
    pptx = _generate_bmode_deck()
    diagram = list(Presentation(pptx).slides)[2]
    pics = _fullbleed_pictures(diagram)
    assert len(pics) >= 1, "다이어그램 슬라이드에 통짜 배경 PICTURE 가 있어야 함(slideBackground 채택)"
    offenders = _overlapping_native_overlays(diagram)
    assert offenders == [], (
        f"다이어그램 통짜 배경 위 겹치는 네이티브 카드/텍스트 {len(offenders)}개 "
        f"(수정2 위반 — build_native_diagram 억제 실패) — {offenders}")


def test_bmode_all_shapes_within_bounds():
    """수정5 검증: 모든 슬라이드의 모든 도형이 경계(0,0,13.333,7.5)in 안."""
    pptx = _generate_bmode_deck()
    for idx, slide in enumerate(list(Presentation(pptx).slides)):
        bad = _out_of_bounds(slide)
        assert bad == [], f"슬라이드 {idx}: 경계 밖 도형 {len(bad)}개 — {bad}"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
