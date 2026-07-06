"""build_native_cover 회귀 테스트 — 도형 생성/편집 가능 여부 측정.

AI는 이미지를 볼 수 없으므로 도형 개수·종류·PICTURE(통짜) 0 여부만 수치 검증한다.
"""
import sys, os
sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "ai_engine"))

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
from native_diagram_pptx import build_native_cover


def _count(slide):
    kinds = {}
    pics = 0
    for sh in slide.shapes:
        t = str(sh.shape_type)
        kinds[t] = kinds.get(t, 0) + 1
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            pics += 1
    return len(slide.shapes._spTree.findall(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}sp')), pics, kinds


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # case 1: 제목+부제+아이브로우+날짜
    s1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    ok1 = build_native_cover(
        s1, prs, title="프로젝트 뎁스별 구조 흐름도",
        subtitle="Bedrock Agents 개발 환경 · Depth 0 → 3 구조 분석",
        eyebrow="Architecture Analysis", date_str="2026-06-22",
        palette=["#4472C4", "#ED7D31", "#A5A5A5"])
    n1, p1, k1 = _count(s1)
    print(f"[case1 title+sub] ok={ok1} shapes(total)={len(s1.shapes)} pictures={p1}")
    assert ok1 is True
    assert p1 == 0, "표지에 통짜 PICTURE가 있으면 안 됨(편집 가능 위배)"
    assert len(s1.shapes) >= 7, f"표지 도형이 너무 적음: {len(s1.shapes)}"

    # case 2: KPI 카드 행
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    ok2 = build_native_cover(
        s2, prs, title="2026 OKR 분기 리뷰",
        subtitle="목표 달성 현황 요약",
        eyebrow="Quarterly Review",
        kpis=[("32%", "매출 성장", "+8%p"), ("95%", "고객 만족", "+3%"),
              ("320", "신규 고객", ""), ("12", "출시 기능", "")],
        palette=["#4472C4", "#ED7D31", "#A5A5A5"])
    n2, p2, k2 = _count(s2)
    print(f"[case2 KPI] ok={ok2} shapes(total)={len(s2.shapes)} pictures={p2}")
    assert ok2 is True
    assert p2 == 0
    assert len(s2.shapes) >= 12, f"KPI 표지 도형 부족: {len(s2.shapes)}"

    # case 3: 빈 제목 → False
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    ok3 = build_native_cover(s3, prs, title="")
    print(f"[case3 empty title] ok={ok3}")
    assert ok3 is False

    out = "/tmp/test_native_cover.pptx"
    prs.save(out)
    print(f"saved → {out}")
    print("ALL PASS")


if __name__ == "__main__":
    main()
