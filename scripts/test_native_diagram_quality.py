"""Regression — 정교화된 네이티브 다이어그램의 기하/시각 품질.

사용자 요구: 렌더러를 훨씬 정교하고 정확하게(젠스파크 수준). 검증 가능한 불변식:
  Q1. flow 가로 카드는 영역 안에 있고 서로 겹치지 않으며 좌→우 순서로 정렬된다.
  Q2. flow에 번호 배지(OVAL)와 화살표 커넥터가 존재한다.
  Q3. org-chart 자식 카드는 겹치지 않고 영역 폭 안에 균등 배치된다.
  Q4. 모든 카드는 슬라이드 region(EMU) 경계 안에 있다(넘침 없음).
  Q5. 텍스트가 긴 노드도 잘리지 않게 카드 폭이 텍스트에 비례해 커진다.

실행: pytest scripts/test_native_diagram_quality.py -q
"""
from __future__ import annotations

import sys
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

pptx = pytest.importorskip("pptx")
nd = pytest.importorskip("native_diagram_pptx")

REGION = (0.6, 1.7, 12.1, 5.2)  # left, top, w, h (inches)


def _blank_slide():
    from pptx import Presentation
    prs = Presentation()
    from pptx.util import Inches
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _shapes(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    boxes, ovals, conns = [], [], []
    for sh in slide.shapes:
        tag = str(getattr(sh._element, "tag", ""))
        if tag.endswith("}cxnSp") or "cxnSp" in tag:
            conns.append(sh)
            continue
        st = sh.shape_type
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if _is_oval(sh):
                ovals.append(sh)
            else:
                boxes.append(sh)
    return boxes, ovals, conns


def _is_oval(sh):
    try:
        return "OVAL" in str(sh.auto_shape_type)
    except Exception:
        return False


def _emu_region():
    from pptx.util import Inches
    l, t, w, h = REGION
    return (Inches(l), Inches(t), Inches(l + w), Inches(t + h))


def test_flow_cards_no_overlap_and_in_bounds():
    """Q1+Q2+Q4 — 가로 흐름 카드 무겹침·정렬·영역내 + 배지/화살표 존재."""
    prs, slide = _blank_slide()
    ok = nd.build_native_diagram(
        slide, "flow", "데이터 수집 -> 전처리 정제 -> 모델 추론 -> 결과 시각화",
        region=REGION)
    assert ok
    boxes, ovals, conns = _shapes(slide)
    # 카드 4개(+배지 4개는 oval로 분리)
    cards = [b for b in boxes]
    assert len(cards) >= 4, f"카드 수 {len(cards)}"
    assert len(ovals) >= 4, f"번호 배지(OVAL) 수 {len(ovals)}"
    assert len(conns) >= 3, f"화살표 커넥터 수 {len(conns)}"

    rl, rt, rr, rb = _emu_region()
    # 카드만 골라 좌→우 정렬 + 무겹침 (배지는 카드 위에 겹칠 수 있어 제외)
    # 카드는 폭이 큰 편 → 너비 상위 4개를 카드로 간주
    cards_sorted = sorted(cards, key=lambda s: s.width, reverse=True)[:4]
    cards_sorted.sort(key=lambda s: s.left)
    for i in range(len(cards_sorted) - 1):
        a, b = cards_sorted[i], cards_sorted[i + 1]
        assert a.left + a.width <= b.left + 9144, (
            f"카드 {i}/{i+1} 가로 겹침: {a.left+a.width} > {b.left}")
    for c in cards_sorted:
        assert c.left >= rl - 9144 and c.left + c.width <= rr + 9144, "카드가 영역 가로 경계 초과"
        assert c.top >= rt - 9144 and c.top + c.height <= rb + 9144, "카드가 영역 세로 경계 초과"


def test_orgchart_children_even_no_overlap():
    """Q3 — 조직도 자식 카드 무겹침 + 영역 폭 내."""
    prs, slide = _blank_slide()
    ok = nd.build_native_diagram(
        slide, "tree", "AI 플랫폼\n  데이터 수집\n  학습 파이프라인\n  서빙 API\n  모니터링",
        region=REGION)
    assert ok
    boxes, ovals, conns = _shapes(slide)
    # 루트1 + 자식4 = 카드 5
    assert len(boxes) >= 5, f"카드 수 {len(boxes)}"
    rl, rt, rr, rb = _emu_region()
    # 자식들(하단, top이 큰 쪽)만: 루트 제외 → top 기준 하위 4개
    children = sorted(boxes, key=lambda s: s.top, reverse=True)[:4]
    children.sort(key=lambda s: s.left)
    for i in range(len(children) - 1):
        a, b = children[i], children[i + 1]
        assert a.left + a.width <= b.left + 9144, f"자식 {i}/{i+1} 겹침"
    for c in children:
        assert c.left >= rl - 9144 and c.left + c.width <= rr + 9144, "자식이 영역 가로 경계 초과"


def test_long_label_widens_card():
    """Q5 — 긴 라벨 카드가 짧은 라벨 카드보다 넓다(텍스트 비례 폭)."""
    prs, slide = _blank_slide()
    ok = nd.build_native_diagram(
        slide, "flow", "A -> 매우 긴 단계 이름 데이터 정제 및 검증 처리 -> B",
        region=REGION)
    assert ok
    boxes, ovals, _ = _shapes(slide)
    cards = sorted(boxes, key=lambda s: s.width, reverse=True)[:3]
    widths = sorted(c.width for c in cards)
    # 가장 넓은 카드가 가장 좁은 카드의 1.3배 이상 (텍스트 비례 반영)
    assert widths[-1] >= widths[0] * 1.3, f"텍스트 비례 폭 미반영: {widths}"


def test_block_has_accent_bars_and_badges():
    """block — 카드 + 좌측 액센트 바 + 번호 배지."""
    prs, slide = _blank_slide()
    ok = nd.build_native_diagram(
        slide, "block", "프론트엔드\n백엔드 API\n데이터베이스", region=REGION)
    assert ok
    boxes, ovals, _ = _shapes(slide)
    # 카드 3 + 액센트 바 3 = autoshape(둥근사각형) 6, 배지 3(oval)
    assert len(boxes) >= 6, f"카드+액센트바 수 {len(boxes)}"
    assert len(ovals) >= 3, f"번호 배지 수 {len(ovals)}"


def test_cards_grid_editable_high_quality():
    """cards — 젠스파크풍 피처 카드 그리드가 편집 가능 네이티브 도형으로 렌더된다(비통짜)."""
    prs, slide = _blank_slide()
    content = ("체계적인 구조: 93개 디렉토리로 구성된 프로젝트\n"
               "효율적인 분류: 5가지 카테고리 자동 분류\n"
               "확장 가능한 아키텍처: 멀티 레벨 구조\n"
               "DevOps 중심: 현대적 관리 접근\n"
               "모듈화된 설계: 유지보수성 극대화")
    ok = nd.build_native_diagram(slide, "cards", content, region=REGION)
    assert ok is True
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    autoshapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    texts = "\n".join(s.text_frame.text for s in slide.shapes if s.has_text_frame)
    pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # 카드(흰 사각형) 5개 + 칩/마크 → autoshape 다수, 통짜 PNG 없음
    assert len(autoshapes) >= 5, f"카드 도형 부족 {len(autoshapes)}"
    assert not pics, "통짜 PNG가 포함됨(편집 불가)"
    # 제목/설명이 편집 가능 텍스트로 들어감
    for t in ("체계적인 구조", "효율적인 분류", "DevOps 중심", "93개 디렉토리"):
        assert t in texts, f"카드 텍스트 누락: {t}"


def test_cards_dispatch_via_build():
    """cards 디스패치 + 빈 입력 안전 폴백."""
    prs, slide = _blank_slide()
    assert nd.build_native_diagram(slide, "cards", "A: 1\nB: 2\nC: 3", region=REGION) is True
    prs2, slide2 = _blank_slide()
    assert nd.build_native_diagram(slide2, "cards", "", region=REGION) is False


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
