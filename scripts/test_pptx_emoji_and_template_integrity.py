"""Regression — (1) 슬라이드/문서 텍스트에서 일반 이모지 제거, (2) 이미지 프롬프트의
전문 아이콘 디렉티브, (3) 템플릿 원본 파일 불변 + 샘플 미유출.

사용자 요구:
  - "일반 이모지 아이콘 절대 사용하지말고, 최신 나노바나나프로급 아이콘 사용."
  - "템플릿 파일 복사본에 새 내용 작성(원본 미변경), 기존 슬라이드에 추가 금지."

수정:
  - `_strip_emoji` + `_normalize_doc_input` 적용: 제목/헤딩/본문/불릿에서 이모지 제거,
    순수 이모지 불릿은 제거, 흐름 화살표(→)·한글은 보존.
  - 이미지 프롬프트(`_build_section_image_prompt`/Vertex)에 "NO EMOJI / premium vector
    line iconography" 디렉티브 추가.
  - 템플릿 생성은 원본 파일을 변경하지 않고(.generated에만 저장) 디자인을 복제해 새 내용 작성.

실행: pytest scripts/test_pptx_emoji_and_template_integrity.py -q
"""
from __future__ import annotations

import os
import sys
import io
import json
import asyncio
import hashlib
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def test_strip_emoji_removes_glyphs_keeps_text():
    assert server._strip_emoji("📁 프로젝트 구조 🚀") == "프로젝트 구조"
    assert server._strip_emoji("백엔드 → 데이터베이스 ✅") == "백엔드 → 데이터베이스"
    assert server._strip_emoji("Frontend 📊 Dashboard") == "Frontend Dashboard"
    # 흐름 화살표는 보존
    assert "→" in server._strip_emoji("A → B ✅")
    # 전부 이모지면 keep_if_empty 기본 True → 원본 보존(제목 소실 방지)
    assert server._strip_emoji("🚀🚀") == "🚀🚀"
    # keep_if_empty=False면 빈 문자열
    assert server._strip_emoji("🚀🚀", keep_if_empty=False) == ""


def test_normalize_strips_emoji_from_slides():
    t, items = server._normalize_doc_input(
        {"title": "📊 보고서",
         "slides": [{"title": "🔧 설정", "bullets": ["✅ 완료", "📁 폴더 → 파일", "🚀🚀"]}]},
        default_kind="slides")
    assert t == "보고서"
    assert items[0]["title"] == "설정"
    # 순수 이모지 불릿은 제거, 나머지는 이모지만 제거(화살표 보존)
    assert items[0]["bullets"] == ["완료", "폴더 → 파일"]


def test_image_prompt_forbids_emoji_and_requests_pro_icons():
    p = server._build_section_image_prompt("flow", "process pipeline steps", "")
    assert "NO EMOJI" in p
    assert "premium vector line iconography" in p
    # 이미지 프롬프트는 영어 디렉티브이므로 한글/이모지 없음
    assert "🚀" not in p


def test_image_prompt_not_truncated_below_directives():
    # [:400]로 늘렸으므로 no_text의 NO EMOJI 디렉티브가 살아있어야 한다
    p = server._build_section_image_prompt("architecture", "system modules layers", "")
    assert "NO EMOJI" in p
    assert len(p) <= 400


_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082")


def _make_template(path, n=6):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    for i in range(n):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = f"ORIG_SAMPLE_{i}"
        if len(s.placeholders) > 1:
            s.placeholders[1].text = f"sample {i}"
        s.shapes.add_picture(io.BytesIO(_PNG), Inches(11), Inches(0.3), Inches(1), Inches(1))
    prs.save(str(path))


def test_template_original_file_unchanged(tmp_path):
    """원본 템플릿 파일은 생성 후 바이트 단위로 변경되지 않는다(복사본에 작성)."""
    from pptx import Presentation
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    tpl = tmp_path / "orig.pptx"
    _make_template(tpl, n=6)
    before = hashlib.md5(tpl.read_bytes()).hexdigest()

    slides = [{"title": "새 A", "bullets": ["x"]}, {"title": "새 B", "bullets": ["y"]}]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "새 덱", "slides": slides, "templatePath": str(tpl)}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        after = hashlib.md5(tpl.read_bytes()).hexdigest()
        assert before == after, "원본 템플릿 파일이 변경됨(복사본 미사용)"
        # 출력은 cover+2=3, 샘플 미유출
        op = Presentation(res["absPath"])
        assert len(op.slides._sldIdLst) == 3
        txt = "\n".join(sh.text_frame.text for sl in op.slides for sh in sl.shapes if sh.has_text_frame)
        assert "ORIG_SAMPLE" not in txt
        assert "새 덱" in txt
    finally:
        ap = res.get("absPath")
        if ap and os.path.isfile(ap):
            try:
                os.remove(ap)
                if os.path.isfile(ap + ".meta.json"):
                    os.remove(ap + ".meta.json")
            except OSError:
                pass


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
