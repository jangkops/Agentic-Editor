"""육안 확인용 데모 — spec: pptx-image-slot-placement-fix.

수정 전(결함)과 수정 후(교정) PPTX 두 벌을 python-pptx 로 직접 합성한다.
D1(풀블리드 중복)·D2(대형 이미지 소형 슬롯)·D3(부분 이미지 슬라이드 밖) 세 결함을
각각 한 슬라이드로 시연하며, before 는 실측 결함 좌표를, after 는
``ai_engine.layout_geometry`` 의 교정 함수를 통과한 좌표를 사용한다.

- before 좌표(실측 결함):
  - D1: 풀블리드 PICTURE 2장 @ (0,0,13.333,7.5)
  - D2: 3840×2160 대형 이미지 @ 0.46in 아이콘 슬롯
  - D3: 부분 이미지 @ (8.11, -1.39, 5.21, 4.17) → top 음수(슬라이드 위로 잘림)
- after 좌표(교정):
  - D1: fullbleed_guard 로 1장만
  - D2: slot_image_fits=False → 콘텐츠 region 으로 승격
  - D3: clamp_into_bounds 로 경계 안

출력: <out_dir>/before.pptx, <out_dir>/after.pptx

네트워크 0(순수 합성). heredoc/stdin 미사용 — 파일로 작성해 실행한다.
"""
from __future__ import annotations

import io
import os
import sys

_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

from PIL import Image, ImageDraw, ImageFont  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt, Emu  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402

import ai_engine.layout_geometry as lg  # noqa: E402

SW, SH = 13.333, 7.5
OUT_DIR = os.path.join(_ROOT, ".generated", "_visual_proof")


def _font(sz):
    for p in (
        "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "/Library/Fonts/Arial.ttf",
    ):
        if os.path.exists(p):
            try:
                return ImageFont.truetype(p, sz)
            except Exception:
                pass
    return ImageFont.load_default()


def _photo_png(path, w, h, label):
    """잘림이 잘 보이도록 그라데이션 + 굵은 테두리 + 큰 라벨/격자 사진 대역 PNG."""
    os.makedirs(os.path.dirname(path), exist_ok=True)
    img = Image.new("RGB", (w, h), (20, 30, 60))
    dr = ImageDraw.Draw(img)
    # 대각 그라데이션
    for y in range(h):
        t = y / max(1, h - 1)
        dr.line([(0, y), (w, y)],
                fill=(int(30 + 180 * t), int(80 + 100 * (1 - t)), int(160 - 80 * t)))
    # 격자(잘림 가시화)
    step = max(40, w // 12)
    for x in range(0, w, step):
        dr.line([(x, 0), (x, h)], fill=(255, 255, 255), width=2)
    for y in range(0, h, step):
        dr.line([(0, y), (w, y)], fill=(255, 255, 255), width=2)
    # 굵은 테두리
    dr.rectangle([3, 3, w - 4, h - 4], outline=(255, 220, 0), width=10)
    # 중앙 라벨 + 모서리 마커(상단이 잘리면 'TOP' 가 사라짐)
    f_big = _font(max(28, h // 8))
    f_sm = _font(max(20, h // 16))
    dr.text((w // 2, h // 2), label, font=f_big, fill=(255, 255, 255), anchor="mm")
    dr.text((w // 2, 24), "▲ TOP EDGE", font=f_sm, fill=(255, 255, 0), anchor="ma")
    dr.text((w // 2, h - 24), "▼ BOTTOM EDGE", font=f_sm, fill=(255, 255, 0), anchor="md")
    img.save(path, "PNG")
    return path


def _label_box(slide, text, *, top=0.08, color=RGBColor(0xF4, 0x47, 0x47)):
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(top), Inches(SW - 0.6), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = color


def _add_pic_clamped(slide, path, rect):
    left, top, wdt, hgt = rect
    return slide.shapes.add_picture(
        path, Inches(left), Inches(top), Inches(wdt), Inches(hgt))


def _blank_layout(prs):
    # 가장 빈 레이아웃(보통 index 6) 사용 — placeholder 간섭 최소화
    return prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]


def _build(fixed: bool):
    """fixed=False → before(결함), True → after(교정)."""
    prs = Presentation()
    prs.slide_width = Emu(int(SW * 914400))
    prs.slide_height = Emu(int(SH * 914400))
    blank = _blank_layout(prs)

    tmp = os.path.join(OUT_DIR, "_assets")
    os.makedirs(tmp, exist_ok=True)
    bg = _photo_png(os.path.join(tmp, "bg4k.png"), 1920, 1080, "BACKGROUND")
    bg2 = _photo_png(os.path.join(tmp, "bg4k_2.png"), 1920, 1080, "BACKGROUND #2")
    big = _photo_png(os.path.join(tmp, "big4k.png"), 1600, 1200, "4K IMAGE")
    illo = _photo_png(os.path.join(tmp, "illo.png"), 900, 720, "ILLUSTRATION")

    tag = "AFTER (교정)" if fixed else "BEFORE (결함)"

    # ── 슬라이드 1: D1 풀블리드 중복 ───────────────────────────────
    s1 = prs.slides.add_slide(blank)
    fb_count = 0
    # 첫 풀블리드(항상 허용)
    if (not fixed) or lg.fullbleed_guard(fb_count):
        _add_pic_clamped(s1, bg, (0, 0, SW, SH)); fb_count += 1
    # 두 번째 풀블리드 후보
    if fixed:
        if lg.fullbleed_guard(fb_count):   # → False, 스킵
            _add_pic_clamped(s1, bg2, (0, 0, SW, SH)); fb_count += 1
    else:
        _add_pic_clamped(s1, bg2, (0, 0, SW, SH)); fb_count += 1  # 결함: 무가드 중복
    _label_box(s1, f"[{tag}] D1 풀블리드 배경 — 임베드된 풀블리드 PICTURE 수 = {fb_count} "
                   f"(기대 ≤ 1)",
               color=RGBColor(0x4E, 0xC9, 0xB0) if fixed else RGBColor(0xF4, 0x47, 0x47))

    # ── 슬라이드 2: D2 대형 이미지 소형 슬롯 ───────────────────────
    s2 = prs.slides.add_slide(blank)
    icon_slot = (1.0, 2.6, 0.46, 0.46)  # 0.46in 아이콘 슬롯
    content_region = (1.5, 1.7, 10.33, 5.2)
    if fixed and not lg.slot_image_fits(icon_slot, 1600, 1200):
        # 콘텐츠 region 으로 승격 + 종횡비 보존 fit
        rect = _fit(content_region, 1600, 1200)
        _add_pic_clamped(s2, big, rect)
        _label_box(s2, f"[{tag}] D2 대형 이미지(1600×1200) → 콘텐츠 영역으로 승격·정상 배치",
                   color=RGBColor(0x4E, 0xC9, 0xB0))
    else:
        # 결함: 0.46in 소형 슬롯에 4K 이미지 욱여넣기
        _add_pic_clamped(s2, big, icon_slot)
        _label_box(s2, "[BEFORE (결함)] D2 대형 이미지(1600×1200)가 0.46in 소형 슬롯에 "
                       "찌그러져 배치 (좌측 작은 점)",
                   color=RGBColor(0xF4, 0x47, 0x47))

    # ── 슬라이드 3: D3 부분 이미지 슬라이드 밖 ─────────────────────
    s3 = prs.slides.add_slide(blank)
    bug_rect = (8.11, -1.39, 5.21, 4.17)  # 실측 결함: top 음수
    if fixed:
        rect = _coords(lg.clamp_into_bounds(bug_rect, (0, 0, SW, SH)))
        _add_pic_clamped(s3, illo, rect)
        _label_box(s3, f"[{tag}] D3 부분 이미지 → clamp 후 경계 안 배치 "
                       f"top={rect[1]:.2f} (≥ 0)",
                   color=RGBColor(0x4E, 0xC9, 0xB0))
    else:
        _add_pic_clamped(s3, illo, bug_rect)  # 결함: 위로 잘림
        _label_box(s3, "[BEFORE (결함)] D3 부분 이미지 top=-1.39in → 상단이 슬라이드 밖으로 "
                       "잘려나감 (▲TOP EDGE 사라짐)",
                   color=RGBColor(0xF4, 0x47, 0x47))

    return prs


def _coords(r):
    if hasattr(r, "left"):
        return (float(r.left), float(r.top), float(r.width), float(r.height))
    s = list(r)
    return (float(s[0]), float(s[1]), float(s[2]), float(s[3]))


def _fit(region, nw, nh):
    return _coords(lg.fit_within(region, nw, nh))


def main():
    os.makedirs(OUT_DIR, exist_ok=True)
    before = _build(fixed=False)
    after = _build(fixed=True)
    bpath = os.path.join(OUT_DIR, "before.pptx")
    apath = os.path.join(OUT_DIR, "after.pptx")
    before.save(bpath)
    after.save(apath)
    print(f"BEFORE: {bpath}")
    print(f"AFTER : {apath}")


if __name__ == "__main__":
    main()
