"""Regression — PPTX 레이아웃 정리: 빈 placeholder 제거 + 네이티브 다이어그램 시 래스터 미중복.

사용자 보고(스크린샷): 생성된 슬라이드가
  (1) "텍스트를 입력하십시오" 빈 placeholder + 점선 테두리가 그대로 노출되고,
  (2) 네이티브 도형(둥근 사각형)과 래스터 AI 이미지가 같은 슬라이드에서 겹쳐
"중구난방 형식"으로 깨졌다. 또한 다이어그램은 편집 가능한 PowerPoint 도형이어야 한다.

근본 원인:
  - 버그 A: `_tool_generate_pptx`에서 native_drawn=True여도 이미지 블록이 별도 `if`라
            래스터 이미지가 추가돼 도형과 겹침.
  - 버그 B: 채우지 않은 본문 placeholder를 `text_frame.clear()`만 해서 PowerPoint가
            프롬프트 "텍스트를 입력하십시오"와 점선을 렌더.

수정:
  - native_drawn이면 슬라이드의 이미지/배경 단계를 전부 건너뛴다(겹침 제거).
  - `_remove_empty_placeholders(slide)`로 빈 placeholder shape 자체를 삭제(제목 보존).

Correctness properties:
  P1. nativeDiagram + imagePrompt가 함께 와도 PPTX에 Picture==0 (도형만, 겹침 없음).
  P2. native_drawn 경로에서 _tool_generate_image는 호출되지 않는다.
  P3. 채우지 않은 본문 placeholder는 슬라이드에서 제거된다(프롬프트 미노출).
  P4. 제목 placeholder는 항상 보존된다.

실행: pytest scripts/test_pptx_layout_native_no_overlap.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def _analyze(pptx_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(pptx_path)
    out = {"autoshapes": 0, "pictures": 0, "slides": []}
    for slide in prs.slides:
        empty_ph = 0
        for ph in slide.placeholders:
            try:
                if not (ph.text_frame.text or "").strip():
                    empty_ph += 1
            except Exception:
                pass
        out["slides"].append({"empty_placeholders": empty_ph})
        for shp in slide.shapes:
            st = shp.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                out["pictures"] += 1
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                out["autoshapes"] += 1
    return out


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_native_diagram_skips_raster_image(tmp_path, monkeypatch):
    """P1+P2 — nativeDiagram + imagePrompt 동시: 래스터 미생성, 도형만 존재."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)

    called = {"n": 0}

    async def _fake_img(tool_input, project_path, aws_profile='', bedrock_user=''):
        # 이 함수가 호출되면(버그) 래스터가 슬라이드에 추가된다.
        called["n"] += 1
        gd = tmp_path / ".generated"
        gd.mkdir(parents=True, exist_ok=True)
        png = gd / "dummy.png"
        # 1x1 PNG (유효 바이트)
        png.write_bytes(bytes.fromhex(
            "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
            "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082"
        ))
        return json.dumps({"path": ".generated/dummy.png"})

    monkeypatch.setattr(server, "_tool_generate_image", _fake_img)

    slides = [{
        "title": "시스템 아키텍처",
        "bullets": ["사용자 -> 게이트웨이 -> 백엔드"],
        "nativeDiagram": {"type": "flow", "content": "사용자 -> 게이트웨이 -> 백엔드"},
        "imagePrompt": "aws architecture diagram, isometric",
    }]
    out = asyncio.run(server._tool_generate_pptx({"title": "T", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        stats = _analyze(res["absPath"])
        assert stats["pictures"] == 0, f"네이티브 도형과 래스터가 겹침(pictures={stats['pictures']})"
        assert called["n"] == 0, "native_drawn인데 이미지 생성이 호출됨(불필요한 래스터 시도)"
        assert stats["autoshapes"] > 0, "편집 가능한 네이티브 도형이 없음"
    finally:
        _cleanup(res)


def test_empty_body_placeholder_removed(tmp_path):
    """P3 — bullets/이미지/다이어그램이 없는 슬라이드의 빈 본문 placeholder는 제거된다."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    # 텍스트가 다이어그램으로 분류되지 않도록 단순 제목만(구조 신호 없음)
    slides = [{"title": "개요"}]
    out = asyncio.run(server._tool_generate_pptx({"title": "T", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        stats = _analyze(res["absPath"])
        # 모든 슬라이드에서 빈 placeholder(프롬프트 "텍스트를 입력하십시오")가 없어야 한다
        for i, sl in enumerate(stats["slides"]):
            assert sl["empty_placeholders"] == 0, (
                f"슬라이드 {i}에 빈 placeholder 잔존({sl['empty_placeholders']}) — 프롬프트 노출")
    finally:
        _cleanup(res)


def test_remove_empty_placeholders_keeps_title():
    """P4 — 헬퍼는 빈 본문 placeholder만 제거하고 제목은 보존한다."""
    from pptx import Presentation
    prs = Presentation()
    # layout 1 = Title and Content (title + body placeholder)
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "제목 유지"
    # 본문 placeholder는 비워둠
    before = len(slide.placeholders._element.getchildren()) if hasattr(slide.placeholders, "_element") else None  # noqa
    server._remove_empty_placeholders(slide)
    # 제목은 남아야 함
    assert slide.shapes.title is not None
    assert (slide.shapes.title.text or "").strip() == "제목 유지"
    # 빈 placeholder는 모두 제거 → 남은 placeholder는 제목뿐(텍스트 보유)
    empties = [ph for ph in slide.placeholders if not (ph.text_frame.text or "").strip()]
    assert empties == [], f"빈 placeholder가 남음: {[ph.placeholder_format.idx for ph in empties]}"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
