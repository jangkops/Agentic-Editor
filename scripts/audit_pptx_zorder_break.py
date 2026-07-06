"""심층 점검: z-order(텍스트가 이미지 아래), 이미지-이미지 겹침, 깨진 이미지 검출.

기존 audit는 텍스트-텍스트 겹침만 봤다. 이번 덱은 슬라이드마다 부분 이미지가 여러 장
올라간 HTML/Vertex 합성 산출물이라 다음을 추가로 검사한다:
  1) z-order 위반: 편집 텍스트 shape가 그것과 겹치는 PICTURE보다 '먼저'(아래) 배치됨
     → 텍스트가 이미지에 가려짐(텍스트가 이미지 아래).
  2) 텍스트↔이미지 겹침: 풀블리드 배경 제외, 부분 이미지가 텍스트와 겹침(임계 낮춤).
  3) 이미지↔이미지 겹침: 부분 이미지끼리 포개짐(레이아웃 깨짐).
  4) 깨진/이상 이미지: 디코드 실패, 0바이트, 극단 종횡비, 슬라이드 밖/초소형.
"""
from __future__ import annotations

import io
import sys

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

try:
    from PIL import Image
    _HAS_PIL = True
except Exception:
    _HAS_PIL = False

EMU = 914400.0
SW, SH = 13.333, 7.5
SLIDE_AREA = SW * SH


def _in(v):
    return round(float(v) / EMU, 3) if v is not None else None


def _rect(sh):
    try:
        l, t, w, h = sh.left, sh.top, sh.width, sh.height
        if None in (l, t, w, h):
            return None
        return (_in(l), _in(t), _in(w), _in(h))
    except Exception:
        return None


def _ov(a, b):
    if not a or not b:
        return 0.0
    ix = max(0.0, min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0]))
    iy = max(0.0, min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1]))
    return ix * iy


def _area(r):
    return max(0.0, r[2]) * max(0.0, r[3]) if r else 0.0


def _fullbleed(r):
    return r and r[0] <= 0.3 and r[1] <= 0.3 and r[2] >= SW * 0.92 and r[3] >= SH * 0.92


def _txt(sh):
    if not getattr(sh, "has_text_frame", False):
        return ""
    try:
        return (sh.text_frame.text or "").strip()
    except Exception:
        return ""


def _img_broken(sh):
    """이미지 깨짐/이상 신호 반환(없으면 None)."""
    try:
        blob = sh.image.blob
    except Exception:
        return "blob 접근 불가"
    if not blob:
        return "0바이트"
    if not _HAS_PIL:
        return None
    try:
        im = Image.open(io.BytesIO(blob))
        im.verify()
        im = Image.open(io.BytesIO(blob))
        w, h = im.size
    except Exception as e:
        return f"디코드 실패({e.__class__.__name__})"
    if w < 8 or h < 8:
        return f"초소형 {w}x{h}"
    ar = w / h if h else 0
    if ar > 20 or ar < 0.05:
        return f"극단 종횡비 {w}x{h}"
    return None


def audit(path):
    prs = Presentation(path)
    print(f"=== 심층 점검(z-order/이미지겹침/깨짐): {path}")
    print(f"슬라이드 {len(prs.slides)}, PIL={_HAS_PIL}\n")
    tot = {"zorder": 0, "txt_img": 0, "img_img": 0, "broken": 0, "offslide": 0}

    for idx, slide in enumerate(prs.slides, 1):
        shapes = list(slide.shapes)
        # z-order = shapes 순서(앞이 아래, 뒤가 위).
        pics, txts = [], []
        for z, sh in enumerate(shapes):
            r = _rect(sh)
            try:
                is_pic = sh.shape_type == MSO_SHAPE_TYPE.PICTURE
            except Exception:
                is_pic = False
            if is_pic:
                pics.append({"z": z, "r": r, "fb": _fullbleed(r), "sh": sh,
                             "broken": _img_broken(sh)})
            else:
                t = _txt(sh)
                if t:
                    txts.append({"z": z, "r": r, "t": t})

        issues = []

        # 1) 깨진 이미지
        for p in pics:
            if p["broken"]:
                tot["broken"] += 1
                issues.append(f"깨진 이미지(z={p['z']}, {p['r']}): {p['broken']}")
            # 슬라이드 밖/초과
            r = p["r"]
            if r and (r[0] < -0.05 or r[1] < -0.05 or r[0] + r[2] > SW + 0.05
                      or r[1] + r[3] > SH + 0.05):
                tot["offslide"] += 1
                issues.append(f"이미지 슬라이드 밖(z={p['z']}, {r})")

        # 2) z-order 위반 + 텍스트↔이미지 겹침 (풀블리드 배경 제외)
        for tx in txts:
            ta = _area(tx["r"])
            if ta <= 0:
                continue
            for p in pics:
                if p["fb"]:
                    continue  # 풀블리드 배경은 의도된 backdrop
                ov = _ov(tx["r"], p["r"])
                if ov <= 0:
                    continue
                pct = ov / ta * 100
                if pct >= 8.0:
                    tot["txt_img"] += 1
                    issues.append(
                        f"텍스트↔이미지 겹침 {pct:.0f}% — '{tx['t'][:18]}'(z={tx['z']}) "
                        f"↔ 이미지(z={p['z']})")
                    # z-order: 텍스트가 이미지보다 먼저(아래) → 가려짐
                    if tx["z"] < p["z"]:
                        tot["zorder"] += 1
                        issues.append(
                            f"  ⚠ z-order 위반 — 텍스트가 이미지 아래(가려짐): "
                            f"'{tx['t'][:18]}'")

        # 3) 이미지↔이미지 겹침 (둘 다 비풀블리드)
        part = [p for p in pics if not p["fb"] and p["r"]]
        for i in range(len(part)):
            for j in range(i + 1, len(part)):
                ov = _ov(part[i]["r"], part[j]["r"])
                amin = min(_area(part[i]["r"]), _area(part[j]["r"]))
                if amin > 0 and ov / amin >= 0.20:
                    tot["img_img"] += 1
                    issues.append(
                        f"이미지↔이미지 겹침 {ov/amin*100:.0f}% "
                        f"(z={part[i]['z']} ↔ z={part[j]['z']})")

        flag = "⚠" if issues else "·"
        print(f"[슬라이드 {idx}] {flag} PIC={len(pics)}(풀블리드 {sum(1 for p in pics if p['fb'])}) "
              f"텍스트={len(txts)}")
        for m in issues:
            print(f"      {m}")

    print("\n=== 요약 ===")
    print(f"깨진/이상 이미지       : {tot['broken']}")
    print(f"슬라이드 밖 이미지     : {tot['offslide']}")
    print(f"텍스트↔이미지 겹침     : {tot['txt_img']}  (그 중 z-order 위반 {tot['zorder']})")
    print(f"이미지↔이미지 겹침     : {tot['img_img']}")


if __name__ == "__main__":
    audit(sys.argv[1] if len(sys.argv) > 1 else "")
