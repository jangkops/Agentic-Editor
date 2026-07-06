"""폴백 격리 단위 테스트 — 요구사항 9 (PPTX 템플릿 처리 실패 격리).

대상: ai_engine/server.py 의 이미 구현된
  - _tool_generate_pptx(tool_input, project_path, aws_profile='', bedrock_user='')  [async]
  - _open_presentation_with_timeout(path, timeout=10)

검증하는 폴백 격리 시나리오:

  1. 템플릿 열기 실패 (요구사항 9.1):
     tool_input["templatePath"]에
       (a) 존재하지 않는 .pptx 경로
       (b) 손상된(.pptx가 아닌 쓰레기 바이트) 파일 경로
     를 주면 → _tool_generate_pptx가 예외를 던지지 않고 무템플릿 경로로 폴백하여
     모든 슬라이드를 포함한 유효 PPTX를 생성한다(path 존재, sizeBytes>0,
     slideCount == 입력 슬라이드 수 + 1(표지)).

  2. style_profile.json 손상 (요구사항 9.2):
     tool_input["styleProfile"]에
       (a) 손상된 JSON 문자열
       (b) 잘못된 색상 값을 가진 dict
       (c) JSON 객체가 아닌 문자열(배열/스칼라)
     를 주면 → SLIDE_DESIGN 기본값으로 폴백하여 생성이 계속되고 유효 PPTX 산출.

  3. 유효성 검증 (요구사항 9.5, 9.6):
     위 실패 케이스에서 생성된 .pptx를 python-pptx Presentation(path)로 다시 열어
     슬라이드 수가 기대대로이고 각 슬라이드에 제목/본문 텍스트가 포함됨을 확인.

  4. 복합 실패 (요구사항 9.6):
     템플릿 열기 실패 + styleProfile 손상을 동시에 주입해도 모든 슬라이드를 포함한
     유효 PPTX 산출이 완료된다.

Bedrock 이미지 호출은 unittest.mock으로 차단한다(슬라이드에 imagePrompt를 두지 않아
실제로는 호출되지 않지만, 네트워크 차단을 위해 방어적으로 patch).

python-pptx 등 필수 의존성 부재 시 skip.

Run:
  ai_engine/.venv/bin/python scripts/test_template_fallback_isolation.py
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile

# repo 루트를 import 경로에 추가 (기존 scripts/ 테스트 컨벤션과 동일).
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

try:
    from pptx import Presentation  # noqa: E402
except Exception as e:  # pragma: no cover - 의존성 부재 시 skip
    print(f"SKIP: python-pptx import 실패 ({e})")
    sys.exit(0)

try:
    from ai_engine.server import _tool_generate_pptx  # noqa: E402
except Exception as e:  # pragma: no cover - server import 실패 시 skip
    print(f"SKIP: ai_engine.server import 실패 ({e})")
    sys.exit(0)


# ---------- 입력 fixture ----------

def _sample_slides():
    """간단한 2개 슬라이드 spec (title/content layout, heading/body 텍스트)."""
    return [
        {
            "title": "첫 번째 슬라이드",
            "layout": "content",
            "bullets": ["본문 항목 A", "본문 항목 B"],
        },
        {
            "title": "두 번째 슬라이드",
            "layout": "two-column",
            "bullets": ["좌측 요점", "우측 요점"],
        },
    ]


def _base_tool_input(extra=None):
    ti = {
        "title": "폴백 격리 테스트 덱",
        "slides": _sample_slides(),
    }
    if extra:
        ti.update(extra)
    return ti


# ---------- 공통 실행/검증 헬퍼 ----------

def _run_generate(tool_input, tmp):
    """_tool_generate_pptx를 동기적으로 실행하고 파싱된 결과 dict를 반환.

    _tool_generate_image를 patch해 Bedrock 호출을 원천 차단한다(방어적).
    """
    from unittest.mock import patch

    async def _no_image(*args, **kwargs):
        return json.dumps({"error": "patched-no-network"})

    with patch("ai_engine.server._tool_generate_image", new=_no_image):
        raw = asyncio.run(_tool_generate_pptx(tool_input, project_path=tmp))
    return json.loads(raw)


def _assert_valid_pptx(parsed, tmp, expected_slide_count, slides):
    """생성 결과가 유효한 PPTX인지 검증 (요구사항 9.5, 9.6).

    - 에러가 없어야 한다 (폴백 격리가 성공해 정상 산출).
    - slideCount == expected_slide_count.
    - path가 .generated/*.pptx 이고 디스크에 존재하며 sizeBytes>0.
    - python-pptx로 다시 열어 슬라이드 수가 일치하고, 각 본문 슬라이드의
      제목·본문 텍스트가 PPTX 내부에 실제로 포함되어 있다.
    """
    assert "error" not in parsed, f"폴백 격리 실패 — 에러 반환됨: {parsed}"

    assert parsed.get("slideCount") == expected_slide_count, (
        f"slideCount 불일치: got {parsed.get('slideCount')}, "
        f"expected {expected_slide_count}; parsed={parsed}"
    )

    rel = parsed.get("path", "")
    assert rel.startswith(".generated/") and rel.endswith(".pptx"), (
        f"예상치 못한 path 형식: {rel}"
    )
    abs_path = os.path.join(tmp, rel)
    assert os.path.isfile(abs_path), f"pptx 파일이 저장되지 않음: {abs_path}"

    assert parsed.get("sizeBytes", 0) > 0, f"sizeBytes가 0 이하: {parsed}"
    assert os.path.getsize(abs_path) > 0, "pptx 파일이 비어 있음"

    # 요구사항 9.5/9.6 — 다시 열어 슬라이드 수와 텍스트 콘텐츠 확인.
    prs = Presentation(abs_path)
    assert len(prs.slides) == expected_slide_count, (
        f"파싱된 슬라이드 수 {len(prs.slides)} != 기대 {expected_slide_count}"
    )

    # 모든 슬라이드의 전체 텍스트를 수집해 입력 제목/본문이 포함됐는지 확인.
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                all_text.append(shape.text_frame.text)
    joined = "\n".join(all_text)

    for sd in slides:
        title_text = sd["title"]
        assert title_text in joined, (
            f"슬라이드 제목 '{title_text}'이(가) 생성 PPTX에 없음.\n전체 텍스트:\n{joined}"
        )
        for bullet in sd.get("bullets", []):
            assert bullet in joined, (
                f"본문 '{bullet}'이(가) 생성 PPTX에 없음.\n전체 텍스트:\n{joined}"
            )

    return abs_path


# ---------- 1. 템플릿 열기 실패 (요구사항 9.1) ----------

def test_template_open_missing_path_falls_back():
    """존재하지 않는 templatePath → 무템플릿 폴백, 유효 PPTX 산출 (9.1, 9.5, 9.6)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        missing = os.path.join(tmp, "does-not-exist.pptx")
        assert not os.path.exists(missing)
        tool_input = _base_tool_input({
            "templatePath": missing,
            "templateId": "tpl-missing",
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
        # 폴백 경로이므로 템플릿 적용 표식(templateId)이 응답에 포함되지 않아야 한다.
        assert "templateId" not in parsed, (
            f"템플릿 열기 실패인데 templateId가 응답에 포함됨: {parsed}"
        )
    print("  9.1 존재하지 않는 templatePath → 무템플릿 폴백 + 유효 PPTX        OK")


def test_template_open_corrupt_file_falls_back():
    """손상된(.pptx 아님) templatePath → 무템플릿 폴백, 유효 PPTX 산출 (9.1, 9.5, 9.6)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        corrupt = os.path.join(tmp, "corrupt-template.pptx")
        # .pptx는 zip 컨테이너인데 임의 바이트를 써서 Presentation() 열기를 깨뜨린다.
        with open(corrupt, "wb") as f:
            f.write(b"this is definitely not a valid pptx / zip container \x00\x01\x02")
        tool_input = _base_tool_input({
            "templatePath": corrupt,
            "templateId": "tpl-corrupt",
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
        assert "templateId" not in parsed, (
            f"손상 템플릿인데 templateId가 응답에 포함됨: {parsed}"
        )
    print("  9.1 손상된 templatePath(.pptx 아님) → 무템플릿 폴백 + 유효 PPTX  OK")


# ---------- 2. style_profile.json 손상 (요구사항 9.2) ----------

def test_corrupt_style_profile_json_string_continues():
    """손상된 styleProfile JSON 문자열 → SLIDE_DESIGN 폴백, 생성 계속 (9.2, 9.5, 9.6)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        tool_input = _base_tool_input({
            # json.loads가 실패하는 깨진 문자열.
            "styleProfile": '{"primaryColor": "#1E1E1E", broken json !!!',
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
    print("  9.2 손상된 styleProfile JSON 문자열 → 기본값 폴백 + 유효 PPTX     OK")


def test_invalid_color_style_profile_dict_continues():
    """잘못된 색상 값을 가진 styleProfile dict → 폴백, 생성 계속 (9.2, 9.5, 9.6).

    _tool_generate_pptx는 styleProfile dict를 그대로 받아 후속 단계에 전달하며,
    여기서 raise하지 않아야 한다. 잘못된 색상이 들어와도 PPTX 산출이 완료되어야 한다.
    """
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        tool_input = _base_tool_input({
            "styleProfile": {
                "primaryColor": "not-a-color",
                "textColor": "#ZZZZZZ",
                "headingFont": "",
                "bodyFont": 12345,  # 비문자열 — 무효
            },
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
    print("  9.2 무효 색상 styleProfile dict → 기본값 폴백 + 유효 PPTX         OK")


def test_non_object_style_profile_string_continues():
    """JSON 객체가 아닌 styleProfile 문자열(배열) → 폴백, 생성 계속 (9.2, 9.5, 9.6)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        tool_input = _base_tool_input({
            # 유효 JSON이지만 dict가 아님 → ValueError 후 None 폴백.
            "styleProfile": "[1, 2, 3]",
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
    print("  9.2 객체 아닌 styleProfile 문자열 → 기본값 폴백 + 유효 PPTX        OK")


# ---------- 3. 복합 실패 (요구사항 9.6) ----------

def test_combined_template_and_profile_failure_still_valid():
    """템플릿 열기 실패 + styleProfile 손상 동시 → 모든 슬라이드 포함 유효 PPTX (9.6)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        corrupt = os.path.join(tmp, "bad.pptx")
        with open(corrupt, "wb") as f:
            f.write(b"\x00 not a pptx \xff")
        tool_input = _base_tool_input({
            "templatePath": corrupt,
            "templateId": "tpl-both-fail",
            "styleProfile": "{ totally broken",
        })
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
        assert "templateId" not in parsed, (
            f"복합 실패인데 templateId가 응답에 포함됨: {parsed}"
        )
    print("  9.6 템플릿+styleProfile 동시 실패 → 모든 슬라이드 포함 유효 PPTX  OK")


# ---------- 4. 기준선(sanity): 폴백 없이 정상 생성 ----------

def test_baseline_no_template_no_profile():
    """templatePath/styleProfile 미전달 시 정상 무템플릿 생성 (기준선 sanity)."""
    slides = _sample_slides()
    with tempfile.TemporaryDirectory() as tmp:
        tool_input = _base_tool_input()
        parsed = _run_generate(tool_input, tmp)
        _assert_valid_pptx(parsed, tmp, len(slides) + 1, slides)
    print("  기준선: 무템플릿/무프로파일 정상 생성                            OK")


def main():
    print("=== 폴백 격리 단위 테스트 (요구사항 9.1, 9.2, 9.5, 9.6) ===")
    test_baseline_no_template_no_profile()
    test_template_open_missing_path_falls_back()
    test_template_open_corrupt_file_falls_back()
    test_corrupt_style_profile_json_string_continues()
    test_invalid_color_style_profile_dict_continues()
    test_non_object_style_profile_string_continues()
    test_combined_template_and_profile_failure_still_valid()
    print("모든 폴백 격리 케이스 통과.")


if __name__ == "__main__":
    main()
