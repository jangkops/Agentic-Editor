#!/usr/bin/env python3
"""Install media generation/editing tools into ai_engine/server.py.

Adds:
- AGENT_TOOLS entries: generate_image, generate_pdf, generate_pptx, edit_image
- _tool_generate_image() — Stability SD3.5 -> Stable Image Core -> Titan Image v2 fallback
- _tool_generate_pdf() — reportlab
- _tool_generate_pptx() — python-pptx
- _tool_edit_image() — Titan v2 inpaint/outpaint -> Nova Canvas fallback
- _execute_tool() routing for the 4 new tools
"""
import os
import re

PATH = os.path.join(os.path.dirname(__file__), '..', 'ai_engine', 'server.py')

with open(PATH, 'r') as f:
    code = f.read()

# ---------------------------------------------------------------------------
# 1. Add new tool entries to AGENT_TOOLS
# ---------------------------------------------------------------------------
new_tools_json = '''            ,
            {
                "toolSpec": {
                    "name": "generate_image",
                    "description": "Generate an image from a text prompt using Bedrock image models. Saves PNG to .generated/. Supports models: Stability SD3.5, Stable Image Core, Amazon Titan Image v2.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "prompt": {"type": "string", "description": "Image description (max 2000 chars)"},
                                "size": {"type": "string", "description": "Image size like 1024x1024, 1024x768, 768x1024", "default": "1024x1024"},
                                "style": {"type": "string", "description": "Optional Stability style preset (photographic, cinematic, anime, etc.)"}
                            },
                            "required": ["prompt"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_pdf",
                    "description": "Generate a multi-page PDF document from structured sections. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Document title (cover page)"},
                                "sections": {
                                    "type": "array",
                                    "description": "List of sections, each with heading and body text",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "heading": {"type": "string"},
                                            "body": {"type": "string"}
                                        }
                                    }
                                }
                            },
                            "required": ["title", "sections"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "generate_pptx",
                    "description": "Generate a PowerPoint presentation from structured slides. Each slide can have a title, bullets, and an optional imagePrompt to auto-generate an image. Saves to .generated/.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "title": {"type": "string", "description": "Cover slide title"},
                                "slides": {
                                    "type": "array",
                                    "description": "List of slides",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "title": {"type": "string"},
                                            "bullets": {"type": "array", "items": {"type": "string"}},
                                            "imagePrompt": {"type": "string", "description": "Optional: auto-generate an image for the slide"},
                                            "layout": {"type": "string", "description": "title | content | two-column", "default": "content"}
                                        }
                                    }
                                }
                            },
                            "required": ["title", "slides"]
                        }
                    }
                }
            },
            {
                "toolSpec": {
                    "name": "edit_image",
                    "description": "Edit an existing PNG/JPEG image using inpaint (mask-based replace) or outpaint (extend canvas). Uses Titan Image v2 with Nova Canvas fallback.",
                    "inputSchema": {
                        "json": {
                            "type": "object",
                            "properties": {
                                "mode": {"type": "string", "description": "inpaint or outpaint"},
                                "image_path": {"type": "string", "description": "Path to the source image"},
                                "prompt": {"type": "string", "description": "Edit instruction (1-512 chars)"},
                                "mask_path": {"type": "string", "description": "Mask image path (inpaint only — white=replace, black=keep)"},
                                "direction": {
                                    "type": "array",
                                    "description": "Outpaint directions: any of left, right, top, bottom",
                                    "items": {"type": "string"}
                                },
                                "extend_pixels": {"type": "integer", "description": "Outpaint extend amount per direction (1-1024)"}
                            },
                            "required": ["mode", "image_path", "prompt"]
                        }
                    }
                }
            }'''

# Find the last tool's closing and inject
# The structure is: "tools": [ {tool1}, {tool2}, ..., {toolN} ]
# Find the closing ] of the search_files tool and add before it
old_close = '''                "required": ["query", "path"]
                            }
                        }
                    }
                }
            }
        }
    ]
}'''

new_close = f'''                "required": ["query", "path"]
                            }}
                        }}
                    }}
                }}
            }}
        }}{new_tools_json}
    ]
}}'''

if old_close not in code:
    print("ERROR: AGENT_TOOLS closing pattern not found")
    print("Searching for alternate...")
    # Try a more lenient match
    m = re.search(r'(\s+"required":\s*\["query",\s*"path"\][\s\S]*?\}\s*\}\s*\}\s*\}\s*)\]\s*\}', code)
    if m:
        # Replace the closing ] with our new tools + ]
        before = code[:m.end(1)]
        after = code[m.end(1):]
        # after starts with "]\n}" — insert new tools before ]
        idx = after.index(']')
        code = before + new_tools_json + after[idx:]
        print("Injected new tools via regex match")
    else:
        print("FATAL: cannot find AGENT_TOOLS structure")
        exit(1)
else:
    code = code.replace(old_close, new_close, 1)
    print("Injected 4 new tools into AGENT_TOOLS")

# ---------------------------------------------------------------------------
# 2. Add tool implementation functions before _execute_tool
# ---------------------------------------------------------------------------
tool_impls = '''

# ===== Media Generation Tools =====

# Image generation model fallback chain
IMAGE_MODELS = [
    "stability.sd3-5-large-v1:0",
    "stability.stable-image-core-v1:1",
    "amazon.titan-image-generator-v2:0",
]
IMAGE_EDIT_MODELS = [
    "amazon.titan-image-generator-v2:0",
    "amazon.nova-canvas-v1:0",
]


async def _tool_generate_image(tool_input: dict, project_path: str) -> str:
    """Generate an image via Bedrock image models with fallback chain.

    Returns JSON string: {path, model, width, height} on success,
    {error, detail?} on failure.
    """
    import time as _t, hashlib, base64
    prompt = (tool_input.get("prompt") or "").strip()
    size = tool_input.get("size", "1024x1024")
    style = tool_input.get("style", "")

    if not prompt:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt is required"})
    if len(prompt) > 2000:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt exceeds 2000 chars"})

    # Parse size
    try:
        w, h = (int(x) for x in size.lower().split("x"))
    except Exception:
        w, h = 1024, 1024

    # Output path
    gen_dir = os.path.join(project_path, ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    ts = str(int(_t.time() * 1000))
    short_hash = hashlib.md5(prompt.encode()).hexdigest()[:4]
    filename = f"image-{ts}-{short_hash}.png"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    aws_profile = os.environ.get("AWS_PROFILE", "bedrock-gw")
    bedrock_user = os.environ.get("BEDROCK_USER", "")
    gw = _get_gw(aws_profile, bedrock_user)

    last_error = ""
    for model_id in IMAGE_MODELS:
        try:
            if model_id.startswith("stability."):
                body = {
                    "prompt": prompt,
                    "mode": "text-to-image",
                    "output_format": "png",
                    "aspect_ratio": "1:1" if w == h else (f"{w}:{h}"),
                }
                if style:
                    body["style_preset"] = style
            elif model_id.startswith("amazon.titan"):
                body = {
                    "textToImageParams": {"text": prompt},
                    "imageGenerationConfig": {
                        "numberOfImages": 1,
                        "width": w,
                        "height": h,
                        "quality": "standard",
                    },
                }
            else:
                body = {"prompt": prompt, "width": w, "height": h}

            result = await gw.invoke_model(model_id, body, timeout=60)

            if "error" in result:
                last_error = f"{model_id}: {result['error'][:200]}"
                continue

            # Extract image bytes
            images = result.get("images", [])
            if not images and isinstance(result, dict):
                # Stability returns "images" with base64 strings
                # Titan returns {"images": ["base64..."]}
                # Some return artifacts
                images = result.get("artifacts", [])
                if images and isinstance(images[0], dict):
                    images = [a.get("base64", "") for a in images]

            if not images:
                last_error = f"{model_id}: no images returned"
                continue

            img_b64 = images[0] if isinstance(images[0], str) else (images[0].get("base64", "") if isinstance(images[0], dict) else "")
            if not img_b64:
                last_error = f"{model_id}: empty image data"
                continue

            img_bytes = base64.b64decode(img_b64)
            with open(output_path, "wb") as f:
                f.write(img_bytes)

            # Get actual dimensions
            try:
                from PIL import Image as _PIL
                with _PIL.open(output_path) as im:
                    aw, ah = im.size
            except Exception:
                aw, ah = w, h

            return json.dumps({
                "path": relative_path,
                "model": model_id,
                "width": aw,
                "height": ah,
                "sizeBytes": len(img_bytes),
            })

        except Exception as e:
            last_error = f"{model_id}: {str(e)[:200]}"
            continue

    return json.dumps({"error": "model-unavailable", "detail": last_error or "all image models failed"})


async def _tool_generate_pdf(tool_input: dict, project_path: str) -> str:
    """Generate a PDF document using reportlab."""
    title = (tool_input.get("title") or "").strip()
    sections = tool_input.get("sections")

    if not title:
        return json.dumps({"error": "invalid-parameter", "detail": "title is required"})
    if not sections or not isinstance(sections, list):
        return json.dumps({"error": "invalid-parameter", "detail": "sections is required (non-empty array)"})

    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.units import cm
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "reportlab", "hint": "pip install reportlab"})

    import time as _t, re as _re
    gen_dir = os.path.join(project_path, ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _re.sub(r"[^a-z0-9]+", "-", title.lower())[:30].strip("-") or "doc"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.pdf"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    try:
        doc = SimpleDocTemplate(output_path, pagesize=A4)
        styles = getSampleStyleSheet()
        story = [Paragraph(title, styles["Title"]), Spacer(1, 1 * cm)]

        for sec in sections:
            heading = sec.get("heading", "") if isinstance(sec, dict) else ""
            body = sec.get("body", "") if isinstance(sec, dict) else str(sec)
            if heading:
                story.append(Paragraph(heading, styles["Heading2"]))
                story.append(Spacer(1, 0.3 * cm))
            if body:
                for para in body.split("\\n"):
                    if para.strip():
                        story.append(Paragraph(para, styles["Normal"]))
                        story.append(Spacer(1, 0.2 * cm))

        doc.build(story)
        size_bytes = os.path.getsize(output_path)
        return json.dumps({
            "path": relative_path,
            "pageCount": len(sections),
            "sizeBytes": size_bytes,
        })
    except Exception as e:
        return json.dumps({"error": "pdf-generation-failed", "detail": str(e)[:200]})


async def _tool_generate_pptx(tool_input: dict, project_path: str) -> str:
    """Generate a PowerPoint presentation using python-pptx."""
    title = (tool_input.get("title") or "").strip()
    slides_data = tool_input.get("slides")

    if not title:
        return json.dumps({"error": "invalid-parameter", "detail": "title is required"})
    if not slides_data or not isinstance(slides_data, list):
        return json.dumps({"error": "invalid-parameter", "detail": "slides is required (non-empty array)"})

    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        return json.dumps({"error": "missing-dep", "lib": "python-pptx", "hint": "pip install python-pptx"})

    import time as _t, re as _re
    gen_dir = os.path.join(project_path, ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
    os.makedirs(gen_dir, exist_ok=True)
    slug = _re.sub(r"[^a-z0-9]+", "-", title.lower())[:30].strip("-") or "deck"
    ts = str(int(_t.time() * 1000))
    filename = f"{slug}-{ts}.pptx"
    output_path = os.path.join(gen_dir, filename)
    relative_path = f".generated/{filename}"

    LAYOUT_MAP = {"title": 0, "content": 1, "two-column": 3}

    try:
        prs = Presentation()
        # Cover slide
        cover_layout = prs.slide_layouts[0]
        cover = prs.slides.add_slide(cover_layout)
        cover.shapes.title.text = title
        if len(cover.placeholders) > 1:
            from datetime import datetime as _dt
            cover.placeholders[1].text = _dt.now().strftime("%Y-%m-%d")

        for i, sd in enumerate(slides_data):
            if not isinstance(sd, dict):
                sd = {"title": str(sd)}
            layout_name = sd.get("layout", "content")
            layout_idx = LAYOUT_MAP.get(layout_name, 1)
            try:
                layout = prs.slide_layouts[layout_idx]
            except IndexError:
                layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = sd.get("title", f"Slide {i + 2}")

            bullets = sd.get("bullets", [])
            body_shape = s.placeholders[1] if len(s.placeholders) > 1 else None
            if body_shape and bullets:
                tf = body_shape.text_frame
                tf.clear()
                for j, bullet in enumerate(bullets):
                    if j == 0:
                        tf.text = str(bullet)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(bullet)

            # Auto-generate image if imagePrompt is set
            img_prompt = sd.get("imagePrompt", "")
            if img_prompt:
                try:
                    img_result_str = await _tool_generate_image(
                        {"prompt": img_prompt, "size": "1024x1024"},
                        project_path,
                    )
                    img_result = json.loads(img_result_str)
                    if "path" in img_result:
                        img_path = os.path.join(project_path, img_result["path"]) if project_path else img_result["path"]
                        if os.path.isfile(img_path):
                            s.shapes.add_picture(img_path, Inches(5), Inches(1.5), width=Inches(4))
                except Exception as e:
                    print(f"[generate_pptx] image gen failed slide {i + 2}: {e}")

        try:
            prs.save(output_path)
        except Exception as save_err:
            return json.dumps({"error": "pptx-generation-failed", "detail": str(save_err)[:200]})

        size_bytes = os.path.getsize(output_path)
        return json.dumps({
            "path": relative_path,
            "slideCount": len(slides_data) + 1,  # +1 for cover
            "sizeBytes": size_bytes,
        })
    except Exception as e:
        return json.dumps({"error": "pptx-generation-failed", "detail": str(e)[:200]})


async def _tool_edit_image(tool_input: dict, project_path: str) -> str:
    """Edit an image using inpaint or outpaint."""
    import time as _t, base64
    mode = tool_input.get("mode", "")
    image_path = tool_input.get("image_path", "")
    prompt = (tool_input.get("prompt") or "").strip()

    if mode not in ("inpaint", "outpaint"):
        return json.dumps({"error": "invalid-mode", "detail": "mode must be inpaint or outpaint"})
    if not image_path:
        return json.dumps({"error": "invalid-parameter", "detail": "image_path is required"})
    if not prompt:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt is required"})
    if len(prompt) > 512:
        return json.dumps({"error": "invalid-parameter", "detail": "prompt exceeds 512 chars"})

    # Resolve path (relative to project_path)
    if not os.path.isabs(image_path) and project_path:
        full_path = os.path.join(project_path, image_path)
    else:
        full_path = image_path
    if not os.path.isfile(full_path):
        return json.dumps({"error": "file-not-found", "detail": f"image not found: {image_path}"})

    # Validate format and size
    try:
        with open(full_path, "rb") as f:
            magic = f.read(8)
        if magic[:8] == b"\\x89PNG\\r\\n\\x1a\\n":
            fmt = "png"
        elif magic[:3] == b"\\xff\\xd8\\xff":
            fmt = "jpeg"
        elif magic[:4] == b"RIFF":
            fmt = "webp"
        else:
            return json.dumps({"error": "invalid-image", "detail": "unsupported format (PNG/JPEG/WEBP only)"})
    except Exception as e:
        return json.dumps({"error": "invalid-image", "detail": str(e)[:200]})

    file_size = os.path.getsize(full_path)
    if file_size > 5 * 1024 * 1024:
        return json.dumps({"error": "invalid-image", "detail": "image exceeds 5MB"})

    # Encode image
    with open(full_path, "rb") as f:
        img_b64 = base64.b64encode(f.read()).decode("ascii")

    # Build request body per mode
    aws_profile = os.environ.get("AWS_PROFILE", "bedrock-gw")
    bedrock_user = os.environ.get("BEDROCK_USER", "")
    gw = _get_gw(aws_profile, bedrock_user)

    last_error = ""

    if mode == "inpaint":
        mask_path = tool_input.get("mask_path", "")
        if not mask_path:
            return json.dumps({"error": "invalid-parameter", "detail": "mask_path required for inpaint"})
        if not os.path.isabs(mask_path) and project_path:
            mask_full = os.path.join(project_path, mask_path)
        else:
            mask_full = mask_path
        if not os.path.isfile(mask_full):
            return json.dumps({"error": "mask-not-found", "detail": f"mask not found: {mask_path}"})
        with open(mask_full, "rb") as f:
            mask_b64 = base64.b64encode(f.read()).decode("ascii")

        for model_id in IMAGE_EDIT_MODELS:
            try:
                if model_id.startswith("amazon.titan"):
                    body = {
                        "taskType": "INPAINTING",
                        "inPaintingParams": {
                            "image": img_b64,
                            "maskImage": mask_b64,
                            "text": prompt,
                        },
                        "imageGenerationConfig": {"numberOfImages": 1, "quality": "standard"},
                    }
                else:  # nova-canvas
                    body = {
                        "taskType": "INPAINTING",
                        "inPaintingParams": {
                            "image": img_b64,
                            "maskImage": mask_b64,
                            "text": prompt,
                        },
                        "imageGenerationConfig": {"numberOfImages": 1},
                    }
                result = await gw.invoke_model(model_id, body, timeout=60)
                if "error" in result:
                    last_error = f"{model_id}: {result['error'][:200]}"
                    continue
                images = result.get("images", [])
                if not images:
                    last_error = f"{model_id}: no images returned"
                    continue
                img_out = images[0] if isinstance(images[0], str) else (images[0].get("base64", "") if isinstance(images[0], dict) else "")
                if not img_out:
                    last_error = f"{model_id}: empty image data"
                    continue
                # Save
                gen_dir = os.path.join(project_path, ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
                os.makedirs(gen_dir, exist_ok=True)
                ts = str(int(_t.time() * 1000))
                filename = f"inpaint-{ts}.png"
                output_path = os.path.join(gen_dir, filename)
                with open(output_path, "wb") as f:
                    f.write(base64.b64decode(img_out))
                try:
                    from PIL import Image as _PIL
                    with _PIL.open(output_path) as im:
                        aw, ah = im.size
                except Exception:
                    aw, ah = 0, 0
                return json.dumps({
                    "path": f".generated/{filename}",
                    "model": model_id,
                    "width": aw,
                    "height": ah,
                    "mode": "inpaint",
                })
            except Exception as e:
                last_error = f"{model_id}: {str(e)[:200]}"
                continue

        return json.dumps({"error": "model-unavailable", "detail": last_error or "all inpaint models failed"})

    # outpaint mode
    direction = tool_input.get("direction", ["right"])
    extend_pixels = int(tool_input.get("extend_pixels", 256))
    if not isinstance(direction, list) or not direction:
        direction = ["right"]
    valid_dirs = {"left", "right", "top", "bottom"}
    direction = [d for d in direction if d in valid_dirs][:4]
    if not direction:
        return json.dumps({"error": "invalid-parameter", "detail": "direction must include left/right/top/bottom"})
    if not (1 <= extend_pixels <= 1024):
        return json.dumps({"error": "invalid-parameter", "detail": "extend_pixels must be 1-1024"})

    for model_id in IMAGE_EDIT_MODELS:
        try:
            if model_id.startswith("amazon.titan"):
                body = {
                    "taskType": "OUTPAINTING",
                    "outPaintingParams": {
                        "image": img_b64,
                        "text": prompt,
                        "outPaintingMode": "DEFAULT",
                    },
                    "imageGenerationConfig": {"numberOfImages": 1, "quality": "standard"},
                }
            else:
                body = {
                    "taskType": "OUTPAINTING",
                    "outPaintingParams": {
                        "image": img_b64,
                        "text": prompt,
                    },
                    "imageGenerationConfig": {"numberOfImages": 1},
                }
            result = await gw.invoke_model(model_id, body, timeout=60)
            if "error" in result:
                last_error = f"{model_id}: {result['error'][:200]}"
                continue
            images = result.get("images", [])
            if not images:
                last_error = f"{model_id}: no images"
                continue
            img_out = images[0] if isinstance(images[0], str) else (images[0].get("base64", "") if isinstance(images[0], dict) else "")
            if not img_out:
                continue
            gen_dir = os.path.join(project_path, ".generated") if project_path else os.path.join(os.getcwd(), ".generated")
            os.makedirs(gen_dir, exist_ok=True)
            ts = str(int(_t.time() * 1000))
            filename = f"outpaint-{ts}.png"
            output_path = os.path.join(gen_dir, filename)
            with open(output_path, "wb") as f:
                f.write(base64.b64decode(img_out))
            try:
                from PIL import Image as _PIL
                with _PIL.open(output_path) as im:
                    aw, ah = im.size
            except Exception:
                aw, ah = 0, 0
            return json.dumps({
                "path": f".generated/{filename}",
                "model": model_id,
                "width": aw,
                "height": ah,
                "mode": "outpaint",
                "direction": direction,
            })
        except Exception as e:
            last_error = f"{model_id}: {str(e)[:200]}"
            continue

    return json.dumps({"error": "model-unavailable", "detail": last_error or "all outpaint models failed"})


'''

# Find _execute_tool definition and insert tool_impls before it
m = re.search(r'\ndef _execute_tool\(', code)
if m:
    code = code[:m.start()] + tool_impls + code[m.start():]
    print("Inserted tool implementations before _execute_tool")
else:
    print("ERROR: _execute_tool not found")
    exit(1)

# ---------------------------------------------------------------------------
# 3. Add routing for new tools in _execute_tool
# ---------------------------------------------------------------------------
# Find the end of search_files branch and add new branches before "else:" or end
old_search = '''        elif tool_name == "search_files":
            query = tool_input["query"]
            path = tool_input["path"]
            if not os.path.isabs(path) and project_path:'''

# Find the entire search_files block + everything until end of try
# Easier: find the routing pattern and add a new generic dispatch after
# Look for the line that closes the if/elif chain — the search_files block
# We add new elif branches before the existing search_files final branch closer

# Strategy: find `elif tool_name == "search_files":` and the try/except wrapping
# Add the new tool dispatches as a new section right after the _REMOTE_TOOLS block
# but use asyncio.run for async functions

# Find: "try:\n        if tool_name == \"read_file\":"
# And add async tool dispatches above it

new_async_dispatch = '''
    # Async media generation tools
    if tool_name in ("generate_image", "generate_pdf", "generate_pptx", "edit_image"):
        try:
            import asyncio as _asyncio
            if tool_name == "generate_image":
                return _asyncio.run(_tool_generate_image(tool_input, project_path))
            if tool_name == "generate_pdf":
                return _asyncio.run(_tool_generate_pdf(tool_input, project_path))
            if tool_name == "generate_pptx":
                return _asyncio.run(_tool_generate_pptx(tool_input, project_path))
            if tool_name == "edit_image":
                return _asyncio.run(_tool_edit_image(tool_input, project_path))
        except Exception as e:
            return json.dumps({"error": "tool-execution-failed", "detail": str(e)[:300]})

'''

# Insert before "try:" of the original tool dispatcher
# Pattern: `    if _BRIDGE_URL and ... _bridge_is_remote():` block ends, then "try:"
# We add right before "try:" that comes before "if tool_name == \"read_file\":"

m2 = re.search(r'(\n    try:\n        if tool_name == "read_file":)', code)
if m2:
    code = code[:m2.start()] + new_async_dispatch + code[m2.start():]
    print("Inserted async dispatch for media tools")
else:
    print("WARNING: could not find tool dispatch insertion point")

with open(PATH, 'w') as f:
    f.write(code)

print(f"Done! File size: {len(code)} chars")
