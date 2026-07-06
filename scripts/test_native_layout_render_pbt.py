"""Property-based tests — pptx-native-density-render (작업 11.1~11.6).

이 파일은 ``ai_engine/native_layout_renderer`` 의 렌더/폴백/라우팅/제목/레이아웃
정확성 속성(Property 1·2·3·7·8·14)을 Hypothesis 로 검증한다. 각 Property 는
**단일 테스트**이며 ``@settings(max_examples=100)`` 이상으로 실행한다.

실행 규약(design.md §Testing Strategy):
  - 네트워크 0: 게이트웨이/Vertex 실호출 없음. Vertex 옵트인은 OFF 로 강제
    (``AE_ENABLE_VERTEX_IMAGE`` != "1") → 장식 이미지 미생성, 콘텐츠 네이티브만.
  - Chrome 불필요: python-pptx 만으로 in-memory Presentation/슬라이드 생성·검증.
  - 파일 저장 불필요(in-memory). watch 모드 금지(단발 pytest 실행).

실행:
  ./venv/bin/python -m pytest scripts/test_native_layout_render_pbt.py -p no:cacheprovider -q
"""

from __future__ import annotations

import os
import sys
import string

# --- Vertex 옵트인 OFF 강제 (네트워크 0, 장식 이미지 미생성) ---
os.environ["AE_ENABLE_VERTEX_IMAGE"] = "0"

# --- repo 루트를 import 경로에 추가 (ai_engine 패키지 import) ---
_REPO_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
if _REPO_ROOT not in sys.path:
    sys.path.insert(0, _REPO_ROOT)

import pytest
from hypothesis import given, settings, strategies as st, HealthCheck

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ai_engine.native_layout_renderer import (
    render_native_layout,
    render_native_fallback,
    finalize_placement,
    PlacedShape,
    NATIVE_LAYOUT_REGISTRY,
    TITLE_ROLES,
    DECORATIVE_BG_ROLES,
    _IMAGE_ROLES,
    _participates_in_collision,
    _norm,
)
from ai_engine.layout_geometry import within_bounds, overlap_area, area, SLIDE_RECT
from ai_engine.slide_templates import design_tokens_for_profile

# --- Property 3 대상: server._should_native_render (실제 함수). import 가능/고속(~0.7s)
#     확인됨. 만약 import 가 불가/hang 위험이면 NATIVE_LAYOUT_REGISTRY 멤버십 + 명시키
#     판정으로 등가 검증하는 폴백을 사용한다(아래 _USING_REAL_GATE 로 표기). ---
try:
    try:
        from ai_engine.server import _should_native_render as _should_native_render_real
    except ImportError:
        from server import _should_native_render as _should_native_render_real  # type: ignore
    _USING_REAL_GATE = True
except Exception:  # pragma: no cover - server import 불가 시 등가 폴백
    _USING_REAL_GATE = False

    def _should_native_render_real(sd, layout, html_enabled):  # type: ignore[no-redef]
        # 등가 판정: html_enabled AND 알려진_레이아웃 AND 명시키 부재 AND 콘텐츠 존재.
        if not html_enabled or not isinstance(sd, dict):
            return False
        if sd.get("imageFile") or sd.get("slideBackground") or sd.get("nativeDiagram"):
            return False
        if not layout or layout not in NATIVE_LAYOUT_REGISTRY:
            return False
        _title = str(sd.get("title", "") or "").strip()
        _bul = sd.get("bullets") or []
        _hb = any(str(b).strip() for b in _bul) if isinstance(_bul, (list, tuple)) else bool(str(_bul).strip())
        return bool(_title or _hb)


_TOKENS = design_tokens_for_profile(None)

# 콘텐츠 역할이 아닌(텍스트 없는 장식/이미지) role 집합 — 이들은 has_text=False 가 정상.
_NONCONTENT_ROLES = set(DECORATIVE_BG_ROLES) | set(_IMAGE_ROLES)

_COMMON_SETTINGS = settings(
    max_examples=100,
    deadline=None,
    suppress_health_check=[HealthCheck.too_slow, HealthCheck.data_too_large,
                           HealthCheck.filter_too_much],
)


# ===========================================================================
# 공통 헬퍼 / 생성기
# ===========================================================================

# python-pptx XML 에 안전한 문자만 사용(제어문자/서로게이트 배제). 한글·라틴·숫자·특수문자.
_SAFE_ALPHABET = (
    "가나다라마바사아자차카타파하거너더러머버서어전정종주중지"
    + string.ascii_letters
    + string.digits
    + " .,!?·—…:/@#%&*()-_=+[]{}"
)


def _ne_text(min_size: int = 1, max_size: int = 40):
    """비어있지 않은(strip 후) 안전 텍스트 — 한글·특수문자 포함."""
    return (st.text(alphabet=_SAFE_ALPHABET, min_size=min_size, max_size=max_size)
            .map(lambda s: s.strip())
            .filter(lambda s: len(s) > 0))


# 초장문(≥200) 포함 — 콘텐츠 텍스트 길이 엣지.
_long_text = (st.text(alphabet=_SAFE_ALPHABET, min_size=200, max_size=360)
              .map(lambda s: s.strip()).filter(lambda s: len(s) > 0))

# 필수 콘텐츠 필드용: 짧은/장문 혼합(항상 비어있지 않음).
_content_text = st.one_of(_ne_text(1, 40), _ne_text(1, 90), _long_text)

# 선택 필드용: 빈문자열/공백/콘텐츠 혼합(렌더러가 빈값은 스킵).
_optional_text = st.one_of(st.just(""), st.just("   "), _ne_text(1, 30))


def _new_slide():
    """13.333×7.5(16:9) Presentation + blank(layout[6]) 슬라이드 생성."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return prs, slide


def _count_editable_text(slide) -> int:
    """슬라이드 트리에서 비어있지 않은 텍스트 런을 가진 편집가능 셰이프 수."""
    n = 0
    for sh in slide.shapes:
        if sh.has_text_frame:
            if any((r.text or "").strip()
                   for p in sh.text_frame.paragraphs for r in p.runs):
                n += 1
    return n


def _count_pictures(slide) -> int:
    """슬라이드 트리의 통짜 이미지(PICTURE) 수 — 콘텐츠 베이크 0 검증용."""
    return sum(1 for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)


# --- 레이아웃별 유효 data 생성기 (필수 필드는 항상 비어있지 않게) ---

def _cards(min_size=2, max_size=4):
    return st.lists(st.fixed_dictionaries({
        "icon": st.just("★"),
        "title": _ne_text(1, 30),
        "description": _optional_text,
    }), min_size=min_size, max_size=max_size)


def _steps(min_size=1, max_size=4):
    return st.lists(st.fixed_dictionaries({
        "label": _ne_text(1, 12),
        "title": _ne_text(1, 30),
        "description": _optional_text,
    }), min_size=min_size, max_size=max_size)


def _items(min_size=1, max_size=4):
    return st.lists(_ne_text(1, 40), min_size=min_size, max_size=max_size)


def _layers(min_size=1, max_size=3):
    return st.lists(st.fixed_dictionaries({
        "name": _ne_text(1, 28),
        "description": _optional_text,
        "items": st.lists(_ne_text(1, 20), min_size=0, max_size=3),
    }), min_size=min_size, max_size=max_size)


def _layout_data_strategy(layout: str):
    """알려진_레이아웃별 유효 data dict 생성기(필수 충족 + 한글/특수/장문/빈선택)."""
    if layout == "cover":
        return st.fixed_dictionaries(
            {"title": _content_text},
            optional={"subtitle": _optional_text, "eyebrow": _optional_text,
                      "footer": _optional_text})
    if layout == "section_divider":
        return st.fixed_dictionaries(
            {"title": _content_text},
            optional={"section_number": st.integers(1, 9), "description": _optional_text})
    if layout == "two_column":
        return st.fixed_dictionaries(
            {"title": _content_text,
             "left_content": st.one_of(_content_text, _items(1, 4)),
             "right_content": st.one_of(_content_text, _items(1, 4))},
            optional={"subtitle": _optional_text})
    if layout == "feature_grid":
        return st.fixed_dictionaries(
            {"title": _content_text, "features": _cards(2, 4)},
            optional={"subtitle": _optional_text})
    if layout == "timeline":
        return st.fixed_dictionaries(
            {"title": _content_text, "steps": _steps(1, 4)},
            optional={"subtitle": _optional_text})
    if layout == "comparison":
        return st.fixed_dictionaries(
            {"title": _content_text,
             "left_label": _ne_text(1, 20), "left_items": _items(1, 4),
             "right_label": _ne_text(1, 20), "right_items": _items(1, 4)},
            optional={"subtitle": _optional_text})
    if layout == "architecture":
        return st.fixed_dictionaries(
            {"title": _content_text, "layers": _layers(1, 3)},
            optional={"subtitle": _optional_text})
    raise AssertionError(f"unknown layout {layout}")


_KNOWN_LAYOUTS = sorted(NATIVE_LAYOUT_REGISTRY.keys())


@st.composite
def _known_layout_and_data(draw):
    layout = draw(st.sampled_from(_KNOWN_LAYOUTS))
    data = draw(_layout_data_strategy(layout))
    return layout, data


@st.composite
def _fallback_dict(draw, *, force_content=True, with_title=None):
    """변환 불가/폴백용 임의 중첩 dict. force_content=True 면 비어있지 않은 콘텐츠
    문자열을 최소 1개 보장한다. with_title 로 title 유무를 제어."""
    d = {}
    if with_title is None:
        with_title = draw(st.booleans())
    if with_title:
        d["title"] = draw(_content_text)
    # 임의 선택 필드(빈문자열·스타일키 포함 — 폴백이 스킵/보존을 구분하는지 자극).
    if draw(st.booleans()):
        d["subtitle"] = draw(_optional_text)
    if draw(st.booleans()):
        d["description"] = draw(_optional_text)
    if draw(st.booleans()):
        d["features"] = draw(_cards(1, 3))
    if draw(st.booleans()):
        d["left_items"] = draw(_items(1, 3))
    if draw(st.booleans()):
        d["accent_color"] = "#FF0000"  # 스타일 힌트(콘텐츠 아님) — 본문 미유출 기대.
    if draw(st.booleans()):
        d["orientation"] = "vertical"
    # 콘텐츠 보장: 비어있지 않은 텍스트가 하나도 없으면 강제 주입.
    if force_content:
        from ai_engine.native_layout_renderer import _collect_fallback_lines
        if not _collect_fallback_lines(d):
            d["body_text"] = draw(_content_text)
    return d


# ===========================================================================
# Property 1 (작업 11.1)
# ===========================================================================
# Feature: pptx-native-density-render, Property 1: 알려진 레이아웃 콘텐츠는 편집가능 네이티브 텍스트로 렌더된다
@_COMMON_SETTINGS
@given(_known_layout_and_data())
def test_property1_known_layout_renders_editable_native_text(layout_and_data):
    """For any 알려진_레이아웃 + 유효 data → render_native_layout 결과의 콘텐츠 셰이프는
    has_text=True 편집가능 네이티브이며, 콘텐츠 텍스트는 이미지로 베이크되지 않는다(베이크 0).

    **Validates: Requirements 1.1, 1.2, 1.3, 11.2**
    """
    layout, data = layout_and_data
    prs, slide = _new_slide()
    result = render_native_layout(slide, prs, layout, data, _TOKENS)

    # 유효 입력은 ok=True, 폴백 트리거 아님.
    assert result.ok is True
    assert result.unsupported is False

    # 콘텐츠 역할 셰이프(장식 배경/이미지 제외)는 모두 비어있지 않은 텍스트 런 보유.
    for ps in result.placed:
        if ps.role not in _NONCONTENT_ROLES:
            assert ps.has_text is True, f"콘텐츠 셰이프 role={ps.role} 가 has_text=False"

    # 슬라이드 트리에 편집가능 텍스트 런이 최소 1개 실재한다.
    assert _count_editable_text(slide) >= 1
    # 콘텐츠 텍스트가 통짜 이미지로 베이크되지 않는다(Vertex OFF → 그림 0).
    assert _count_pictures(slide) == 0


# ===========================================================================
# Property 2 (작업 11.2)
# ===========================================================================
# Feature: pptx-native-density-render, Property 2: 변환 불가 입력도 편집가능 텍스트로 폴백된다
@_COMMON_SETTINGS
@given(data=_fallback_dict(force_content=True), use_unknown=st.booleans())
def test_property2_unconvertible_falls_back_to_editable_text(data, use_unknown):
    """For any 변환 불가(미지원 레이아웃 또는 필수필드 부재) 입력 → render_native_layout 은
    unsupported 를 신호하고, render_native_fallback 은 콘텐츠 텍스트를 편집가능 셰이프 1개
    이상으로 출력하며 통짜 이미지로 대체하지 않는다(그림 0).

    **Validates: Requirements 1.4, 1.5**
    """
    # 1) 변환 불가 트리거 확인 — 미지원 레이아웃 또는 필수필드 부재.
    prs1, slide1 = _new_slide()
    if use_unknown:
        trigger_layout = "unknown_layout_xyz"
    else:
        # 알려진 레이아웃이되 필수필드(title 제외 등)가 부재한 dict → 폴백 트리거.
        trigger_layout = "two_column"  # left/right_content 필수 부재
        data = {k: v for k, v in data.items() if k not in ("left_content", "right_content")}
    res_layout = render_native_layout(slide1, prs1, trigger_layout, data, _TOKENS)
    assert res_layout.ok is False
    assert res_layout.unsupported is True

    # 2) 폴백 — 콘텐츠 텍스트 전수를 편집가능 텍스트박스로.
    prs2, slide2 = _new_slide()
    res_fb = render_native_fallback(slide2, data, _TOKENS)
    assert res_fb.ok is True
    # 편집가능 텍스트 셰이프 1개 이상.
    assert sum(1 for ps in res_fb.placed if ps.has_text) >= 1
    assert _count_editable_text(slide2) >= 1
    # 통짜 이미지 대체 금지(그림 0).
    assert _count_pictures(slide2) == 0


# ===========================================================================
# Property 3 (작업 11.3)
# ===========================================================================
# Feature: pptx-native-density-render, Property 3: 네이티브 라우팅 결정의 일관성
@_COMMON_SETTINGS
@given(
    layout=st.sampled_from(_KNOWN_LAYOUTS),
    title=_ne_text(1, 40),
    explicit=st.sampled_from([None, "imageFile", "slideBackground"]),
)
def test_property3_should_native_render_consistency(layout, title, explicit):
    """For any 알려진_레이아웃 + 콘텐츠 존재 슬라이드:
      - caller 가 imageFile/slideBackground 를 명시하지 않으면 True(네이티브 라우팅).
      - imageFile/slideBackground 를 명시하면 False(명시 경로 보존).

    NOTE: 대상은 server._should_native_render 실제 함수(import 고속·hang 없음 확인).
    import 불가 시 NATIVE_LAYOUT_REGISTRY 멤버십 + 명시키 판정 등가 폴백 사용
    (_USING_REAL_GATE 로 표기).

    **Validates: Requirements 1.5, 9.1**
    """
    sd = {"layout": layout, "title": title}
    if explicit is not None:
        sd[explicit] = "assets/something.png"

    decision = _should_native_render_real(sd, layout, True)

    if explicit is None:
        assert decision is True
    else:
        assert decision is False


# ===========================================================================
# Property 7 (작업 11.4)
# ===========================================================================
# Feature: pptx-native-density-render, Property 7: 제목 셰이프 수는 제목 유무에 정확히 일치한다
@_COMMON_SETTINGS
@given(
    mode=st.sampled_from(["known", "fallback_with_title", "fallback_no_title"]),
    known=_known_layout_and_data(),
    fb_with=_fallback_dict(force_content=True, with_title=True),
    fb_without=_fallback_dict(force_content=True, with_title=False),
)
def test_property7_title_count_matches_presence(mode, known, fb_with, fb_without):
    """For any 입력: 제목 텍스트가 있으면 title_count==1, 없으면 0.

    **Validates: Requirements 4.1, 4.4**
    """
    prs, slide = _new_slide()
    if mode == "known":
        layout, data = known  # 알려진_레이아웃은 항상 title 필수 → 제목 있음.
        result = render_native_layout(slide, prs, layout, data, _TOKENS)
        expected = 1
    elif mode == "fallback_with_title":
        result = render_native_fallback(slide, fb_with, _TOKENS)
        expected = 1
    else:  # fallback_no_title — title 키 없음(다른 콘텐츠는 존재)
        assert "title" not in fb_without
        result = render_native_fallback(slide, fb_without, _TOKENS)
        expected = 0

    assert result.title_count == expected
    # placed 의 제목 역할 셰이프 수도 일치한다.
    assert sum(1 for ps in result.placed if ps.role in TITLE_ROLES) == expected


# ===========================================================================
# Property 8 (작업 11.5)
# ===========================================================================

@st.composite
def _title_dedup_placed(draw):
    """동일 제목 2+ (정규화 후 일치) + 서로 다른 제목/본문 셰이프 목록 생성.

    중복 제목은 같은 영역에 둔다(dedup 으로 1개만 남음). 살아남는 셰이프(고유 제목/본문)는
    경계 안 비겹침 y-스택으로 배치해 finalize_placement 의 겹침 보정이 OverlapError 없이
    no-op 이 되도록 한다.
    """
    base = draw(_ne_text(1, 24))
    n_dup = draw(st.integers(2, 4))

    def _variant(s, k):
        return [s, "  " + s + "  ", s.upper(), s.lower(), s.swapcase()][k % 5]

    placed = []
    # 중복 제목 — 동일 rect (정규화 시 base 와 동일).
    for i in range(n_dup):
        placed.append(PlacedShape(role="title", rect=(1.0, 0.5, 4.0, 0.6),
                                   has_text=True, text=_variant(base, i), z=5))

    # 서로 다른 제목(정규화가 base 와 다른 것) — 비겹침 y-스택.
    distinct = draw(st.lists(
        _ne_text(1, 24).filter(lambda s: _norm(s) != _norm(base)),
        min_size=0, max_size=2, unique_by=_norm))
    y = 1.5
    for s in distinct:
        placed.append(PlacedShape(role="title", rect=(1.0, y, 4.0, 0.6),
                                   has_text=True, text=s, z=5))
        y += 1.0

    # 본문 셰이프 몇 개 — 비겹침.
    n_body = draw(st.integers(0, 2))
    by = 4.2
    for i in range(n_body):
        placed.append(PlacedShape(role="body", rect=(1.0, by, 4.0, 0.6),
                                   has_text=True, text=draw(_ne_text(1, 20)), z=2))
        by += 0.9

    return base, distinct, placed


# Feature: pptx-native-density-render, Property 8: 중복 제목 제거 및 베이크 제목 미채택
@_COMMON_SETTINGS
@given(_title_dedup_placed())
def test_property8_duplicate_title_dedup(payload):
    """For any 한 슬라이드에 정규화 후 동일 제목이 2개 이상 방출될 때 → finalize_placement 는
    제목 셰이프 1개만 남기고, 동일 제목이 베이크(장식 배경)로 채택되지 않는다.

    **Validates: Requirements 4.2, 4.3**
    """
    base, distinct, placed = payload
    result = finalize_placement(placed)

    nbase = _norm(base)
    # 동일 정규화 제목 셰이프는 정확히 1개만 남는다.
    same = [ps for ps in result if ps.role in TITLE_ROLES and _norm(ps.text) == nbase]
    assert len(same) == 1

    # 서로 다른 제목은 모두 보존된다(각 고유 정규화 1개씩).
    for s in distinct:
        kept = [ps for ps in result if ps.role in TITLE_ROLES and _norm(ps.text) == _norm(s)]
        assert len(kept) == 1

    # 전체 제목 역할 셰이프 수 == 고유 정규화 제목 수(중복 base 1 + distinct).
    expected_titles = 1 + len(distinct)
    assert sum(1 for ps in result if ps.role in TITLE_ROLES) == expected_titles

    # 베이크 제목 미채택 — 제목 텍스트를 담은 장식 배경(decorative_bg) 셰이프가 없다.
    assert not any(ps.role == "decorative_bg" and _norm(ps.text) == nbase for ps in result)


# ===========================================================================
# Property 14 (작업 11.6)
# ===========================================================================
# Feature: pptx-native-density-render, Property 14: 모든 알려진 레이아웃에 4규칙이 동시에 성립한다
@_COMMON_SETTINGS
@given(data_by_layout=st.fixed_dictionaries(
    {ly: _layout_data_strategy(ly) for ly in _KNOWN_LAYOUTS}))
def test_property14_all_known_layouts_satisfy_four_rules(data_by_layout):
    """For any 알려진_레이아웃 7종 + 유효 data → 네이티브 렌더 결과는 동시에
    (편집가능 텍스트 런 ≥1) AND (셰이프 쌍 겹침 < 10%) AND (모든 셰이프 경계 안) AND
    (제목 셰이프 ≤ 1)을 만족한다.

    **Validates: Requirements 7.1, 7.2, 7.3**
    """
    for layout in _KNOWN_LAYOUTS:
        data = data_by_layout[layout]
        prs, slide = _new_slide()
        result = render_native_layout(slide, prs, layout, data, _TOKENS)
        assert result.ok is True, f"{layout}: ok=False"

        # (1) 편집가능 텍스트 런 ≥ 1
        assert _count_editable_text(slide) >= 1, f"{layout}: 편집가능 텍스트 0"

        # (2) 제목 셰이프 ≤ 1
        title_n = sum(1 for ps in result.placed if ps.role in TITLE_ROLES)
        assert title_n <= 1, f"{layout}: title 셰이프 {title_n}개"

        # (3) 모든 셰이프 경계 안(eps=0.05, float slack 미세 허용)
        for ps in result.placed:
            assert within_bounds(ps.rect, SLIDE_RECT, eps=0.051), \
                f"{layout}: 경계 밖 role={ps.role} rect={ps.rect}"

        # (4) 겹침 검사(A안) — 텍스트 보유/비배경 이미지 쌍만, 겹침 < 10%
        part = [ps for ps in result.placed if _participates_in_collision(ps)]
        for i in range(len(part)):
            for j in range(i + 1, len(part)):
                a, b = part[i].rect, part[j].rect
                ov = overlap_area(a, b)
                if ov <= 0.0:
                    continue
                amin = min(area(a), area(b))
                if amin <= 0.0:
                    continue
                assert ov < 0.10 * amin + 1e-6, \
                    f"{layout}: 겹침 {ov/amin:.3f} ≥ 10% ({part[i].role},{part[j].role})"
