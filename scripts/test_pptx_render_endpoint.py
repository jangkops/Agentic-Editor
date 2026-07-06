"""Regression — /api/media/pptx-render 가 생성된 .pptx를 미리보기용으로 파싱한다.

사용자 보고: 앱에서 'POST /api/media/pptx-render → 404'가 발생해 생성물 미리보기가
안 되고 품질 확인이 불가. 근본 원인: 서버에 해당 라우트가 미구현(프론트 pptx-viewer.js만
존재). 수정: python-pptx로 슬라이드 제목/본문 + 임베드 이미지(base64 data URL)를 반환.

Correctness properties:
  P1. 슬라이드별 title/bullets/images 키를 가진 리스트를 반환한다.
  P2. 임베드 Picture가 data URL(base64)로 인라인된다(미리보기에 시각 표시 가능).
  P3. 존재하지 않는 path는 404 에러를 반환한다(서버 크래시 없음).

실행: pytest scripts/test_pptx_render_endpoint.py -q
"""
from __future__ import annotations

import io
import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")

_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000154a24f9f0000000049454e44ae426082")


class _FakeReq:
    def __init__(self, payload):
        self._p = payload

    async def json(self):
        return self._p


def _make_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "프로젝트 흐름도"
    if len(s.placeholders) > 1:
        s.placeholders[1].text = "단계 A\n단계 B"
    s.shapes.add_picture(io.BytesIO(_PNG), Inches(1), Inches(2), Inches(4), Inches(3))
    prs.save(str(path))


def _body(resp):
    # JSONResponse → bytes 본문 파싱
    raw = resp.body
    if isinstance(raw, (bytes, bytearray)):
        return json.loads(raw.decode("utf-8"))
    return json.loads(raw)


def test_pptx_render_parses_slides_and_images(tmp_path):
    """P1+P2 — 제목/본문/이미지(dataURL) 반환."""
    pptx_path = tmp_path / "deck.pptx"
    _make_pptx(pptx_path)
    resp = asyncio.run(server.api_media_pptx_render(_FakeReq({"path": str(pptx_path)})))
    data = _body(resp)
    assert "slides" in data and len(data["slides"]) == 1
    sl = data["slides"][0]
    assert sl["title"] == "프로젝트 흐름도"
    assert "단계 A" in sl["bullets"] and "단계 B" in sl["bullets"]
    assert sl["images"], "임베드 이미지가 반환되지 않음"
    assert sl["images"][0].startswith("data:image/"), "이미지가 data URL이 아님"


def test_pptx_render_missing_file_404(tmp_path):
    """P3 — 없는 경로는 404."""
    resp = asyncio.run(server.api_media_pptx_render(_FakeReq({"path": str(tmp_path / "nope.pptx")})))
    assert resp.status_code == 404


def test_pptx_render_requires_path():
    """path 누락 시 400."""
    resp = asyncio.run(server.api_media_pptx_render(_FakeReq({})))
    assert resp.status_code == 400


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
