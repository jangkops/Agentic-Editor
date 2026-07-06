#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""무템플릿 PPTX 공통 디자인 베이스라인 회귀 테스트.

내용/다이어그램 종류와 무관하게 모든 슬라이드에 상단 액센트 바 + 제목 스타일 +
하단 푸터(문서명·페이지)가 자동 적용되는지 검증. 편집 가능(통짜 PICTURE 아님) 유지.
"""
import os
import sys
import asyncio
import json
import tempfile

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

import pytest


def _gen(tmp, tool_input, monkeypatch):
    # monkeypatch.setenv는 테스트 종료 시 자동 복원 → 프로세스 환경 누수 방지
    # (이전: os.environ 직접 설정으로 AE_DISABLE_HTML_SLIDES가 누수돼 HTML 슬라이드
    #  테스트를 깨뜨렸음 — 테스트 격리 결함).
    monkeypatch.setenv("AE_GENERATED_ROOT", tmp)
    monkeypatch.setenv("AE_ENABLE_HTML_SLIDES", "0")
    monkeypatch.setenv("AE_DISABLE_HTML_SLIDES", "1")
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "0")
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "1")
    monkeypatch.setenv("AE_PPTX_TOC", "0")
    # Vertex 배경(네트워크 기능)은 구조 불변식 검증과 무관 — 결정성 위해 비활성.
    monkeypatch.setenv("AE_DISABLE_VERTEX_IMAGE", "1")
    monkeypatch.setenv("AE_DISABLE_VERTEX_HERO", "1")
    monkeypatch.setenv("AE_DISABLE_VERTEX_BODY_BG", "1")
    from ai_engine import server
    res = asyncio.run(server._tool_generate_pptx(tool_input, project_path=tmp))
    obj = json.loads(res)
    path = obj.get("absPath") or obj.get("path")
    cand = path if os.path.isabs(path) else os.path.join(tmp, path)
    if not os.path.isfile(cand):
        cand = os.path.join(tmp, ".generated", os.path.basename(path))
    assert os.path.isfile(cand), f"pptx not found: {path}"
    return cand


def test_universal_design_on_plain_slides(monkeypatch):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    tmp = tempfile.mkdtemp(prefix="ae_design_")
    cand = _gen(tmp, {
        "title": "사업 전략 보고",
        "slides": [
            {"title": "시장 개요", "bullets": ["성장세 지속", "경쟁 심화", "니즈 다변화"]},
            {"title": "우리의 강점", "bullets": ["기술력", "고객 기반", "빠른 실행력"]},
        ],
    }, monkeypatch)
    prs = Presentation(cand)
    assert len(prs.slides) == 3  # 표지 + 2

    for i, slide in enumerate(prs.slides):
        autos = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        # 상단 액센트 바 + 틱 → 최소 2개 도형
        assert len(autos) >= 2, f"slide {i} 공통 디자인 도형 부족: {len(autos)}"
        # 편집 가능 불변식 — 통짜(풀블리드) PICTURE 없음. 작은 장식 아이콘 칩은 허용.
        fullbleed = [p for p in pics
                     if (p.width or 0) >= prs.slide_width * 0.9
                     and (p.height or 0) >= prs.slide_height * 0.9]
        assert not fullbleed, f"slide {i} 통짜 풀블리드 PICTURE 발견"

    # 본문 슬라이드: 페이지 번호 "i / n" 텍스트 존재
    for i, slide in enumerate(prs.slides):
        if i == 0:
            continue
        texts = [s.text_frame.text for s in slide.shapes
                 if s.has_text_frame and (s.text_frame.text or "").strip()]
        joined = " ".join(texts)
        assert f"{i + 1} / 3" in joined, f"slide {i} 페이지 번호 누락: {texts}"


def test_universal_design_does_not_break_diagram_slides(monkeypatch):
    """다이어그램(flow) 슬라이드에도 공통 디자인이 함께 적용되고 깨지지 않는다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    tmp = tempfile.mkdtemp(prefix="ae_design2_")
    cand = _gen(tmp, {
        "title": "프로세스 안내",
        "slides": [
            {"title": "진행 흐름", "bullets": ["계획", "개발", "테스트", "출시"]},
        ],
    }, monkeypatch)
    prs = Presentation(cand)
    # 표지 + 1
    assert len(prs.slides) == 2
    last = prs.slides[1]
    autos = [s for s in last.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    pics = [s for s in last.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    # 다이어그램 도형 + 공통 디자인 도형
    assert len(autos) >= 4
    # 통짜(풀블리드) PICTURE 없음. 작은 장식 아이콘 칩은 허용.
    fullbleed = [p for p in pics
                 if (p.width or 0) >= prs.slide_width * 0.9
                 and (p.height or 0) >= prs.slide_height * 0.9]
    assert not fullbleed


if __name__ == "__main__":
    sys.exit(pytest.main([__file__, "-v", "-p", "no:cacheprovider"]))
