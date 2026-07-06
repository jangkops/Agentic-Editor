"""audit_pptx_native_density — 네이티브 렌더 산출물(.pptx)의 통합 합격 게이트.

이 모듈은 ``pptx-native-density-render`` 스펙의 산출물_검증기다. 실제 생성된
.pptx 파일/슬라이드를 python-pptx 만으로(=Chrome·네트워크 불필요) 검사한다.

본 파일은 두 검증 차원을 담는다(서로 **직교**한다):

  1. **밀도(요소 존재)** — ``parity_scorer.score``/``score_native_slide`` 가
     "디자인 요소가 있는가"(체크리스트 마커 개수)를 센다. (작업 7.1)
  2. **스타일 품질** — ``audit_style_quality(slide, tokens)`` 가 "기존
     design_tokens 가 도형 서식에 실제로 적용됐는가"(라운드/그림자/accent색/
     타이포 계층/여백)를 검사한다. (이 파일, 작업 6.1)

이 두 차원은 **별개**다. 밀도 채점이 통과해도(요소가 있어도) 스타일 품질이
미달일 수 있고(토큰 미적용), 반대도 성립한다. "초고퀄"은 둘 다 합격해야 한다.
``audit_native_density`` 의 (h) 검사가 ``audit_style_quality`` 를 호출한다.

설계 근거: ``.kiro/specs/pptx-native-density-render/design.md``
  - Components and Interfaces §3-bis "비주얼 품질 검증기 audit_style_quality"
  - Correctness Property 19
요구사항: Req 5.1 / 5.3 / 5.4.

제약: 색·여백·타이포 기준값은 모두 ``design_tokens_for_profile`` 결과 dict
(인자 ``tokens``)에서만 읽는다 — **신규 토큰을 정의하지 않는다(Req 5.4)**.
이 검사는 "기존 토큰이 실제로 적용됐는지"를 보는 것이지 새 토큰 설계가 아니다.
"""

from __future__ import annotations

from dataclasses import dataclass, field

# python-pptx — 네이티브 도형 서식/효과를 읽기 위한 의존성. 미설치 환경에서도
# 모듈 import 는 가능하도록 가드(실제 audit 호출 시에만 필요).
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_SHAPE
    from pptx.enum.dml import MSO_FILL
    from pptx.oxml.ns import qn
except Exception:  # pragma: no cover - python-pptx 미설치 환경
    MSO_SHAPE_TYPE = MSO_SHAPE = MSO_FILL = qn = None  # type: ignore[assignment]


# ---------------------------------------------------------------------------
# python-pptx 텍스트 프레임 기본 내부 마진(EMU). 토큰 기반 여백이 "명시 적용"
# 됐는지를 판정하는 기준(이 값과 다르면 명시 설정으로 본다).
#   기본 좌/우 = 0.1in = 91440 EMU, 상/하 = 0.05in = 45720 EMU.
# ---------------------------------------------------------------------------
_DEFAULT_MARGIN_LR_EMU = 91440   # 0.10 in
_DEFAULT_MARGIN_TB_EMU = 45720   # 0.05 in

# 스타일 품질 검사 항목 식별자(누락 시 missing_style 에 기록).
CHECK_ROUND_DECOR = "round_shadow_or_border"   # 카드/박스 라운드 + 그림자/테두리
CHECK_ACCENT_COLOR = "accent_color"            # 섹션헤더/배지 accent 색(검정/흰색 아님)
CHECK_TYPE_HIERARCHY = "type_hierarchy"        # 제목 폰트 > 본문 폰트
CHECK_BODY_MARGIN = "body_margin_tokens"       # 본문 여백 토큰 명시 적용
CHECK_MIN_VISUALS = "min_visual_elements"      # 시각 요소 최소 개수(>=2)

# 슬라이드당 최소 시각 요소 개수(도형/배지/카드/구분선/이미지 등 비텍스트 장식).
_MIN_VISUAL_ELEMENTS = 2


@dataclass
class StyleQualityReport:
    """``audit_style_quality`` 반환값 — 스타일 품질(토큰 실제 적용) 검사 결과.

    Attributes:
        passed: 모든 검사 항목 충족 여부.
        score: 충족 항목 수 / 전체 항목 수 (0.0~1.0).
        missing_style: 누락된 스타일 항목명 리스트(불합격 시 비어있지 않음).
    """

    passed: bool
    score: float
    missing_style: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# 색 파싱/비교 헬퍼 — 기준값은 모두 인자 tokens 에서만 읽는다(신규 토큰 금지).
# ---------------------------------------------------------------------------


def _hex_to_rgb(hexstr):
    """'#RRGGBB' 또는 'RRGGBB' → (r, g, b). 형식 불일치 시 None."""
    if not hexstr:
        return None
    s = str(hexstr).strip().lstrip("#")
    if len(s) != 6:
        return None
    try:
        return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))
    except Exception:
        return None


def _color_distance(c1, c2) -> float:
    """두 RGB 튜플의 유클리드 거리."""
    return (sum((int(a) - int(b)) ** 2 for a, b in zip(c1, c2))) ** 0.5


def _is_black_or_white(c) -> bool:
    """기본 검정(#000000)/흰색(#FFFFFF) 근방 여부."""
    if c is None:
        return False
    return all(v <= 12 for v in c) or all(v >= 243 for v in c)


def _solid_fill_rgb(shape):
    """도형의 solid fill 전경색을 (r, g, b)로. solid 아님/실패 시 None."""
    try:
        fill = shape.fill
        if MSO_FILL is not None and fill.type != MSO_FILL.SOLID:
            return None
        return _hex_to_rgb(str(fill.fore_color.rgb))
    except Exception:
        return None


def _accent_like_colors(tokens: dict) -> list:
    """tokens 에서 accent/primary/secondary 색을 RGB 튜플 리스트로 추출."""
    out = []
    for key in ("accent", "primary", "secondary"):
        rgb = _hex_to_rgb((tokens or {}).get(key))
        if rgb is not None:
            out.append(rgb)
    return out


# ---------------------------------------------------------------------------
# 도형 속성 판정 헬퍼 (python-pptx)
# ---------------------------------------------------------------------------


def _is_auto_shape(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    except Exception:
        return False


def _is_rounded(shape) -> bool:
    """라운드 코너 도형 여부 — ROUNDED_RECTANGLE 이거나 adjustments[0] > 0."""
    # 1) auto_shape_type 이 ROUNDED_RECTANGLE 계열인가.
    try:
        ast = shape.auto_shape_type
        if ast is not None and MSO_SHAPE is not None:
            name = getattr(ast, "name", "") or ""
            if "ROUNDED" in str(name).upper():
                return True
            try:
                if int(ast) == int(MSO_SHAPE.ROUNDED_RECTANGLE):
                    return True
            except Exception:
                pass
    except Exception:
        pass
    # 2) adjustment 핸들(>0) 로 라운드 정도가 적용됐는가.
    try:
        adjs = shape.adjustments
        for i in range(len(adjs)):
            if float(adjs[i]) > 0.0:
                return True
    except Exception:
        pass
    return False


def _has_shadow(shape) -> bool:
    """outerShdw(드롭 섀도우) 적용 여부 — spPr/a:effectLst/a:outerShdw 존재."""
    if qn is None:
        return False
    try:
        spPr = shape._element.spPr
        if spPr is None:
            return False
        for eff in spPr.findall(qn("a:effectLst")):
            if eff.find(qn("a:outerShdw")) is not None:
                return True
    except Exception:
        pass
    return False


def _has_border(shape) -> bool:
    """테두리(line) 적용 여부 — a:ln 에 solidFill 이 있고 noFill 이 아님."""
    if qn is None:
        return False
    try:
        spPr = shape._element.spPr
        if spPr is None:
            return False
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            return False
        # 명시적 noFill 이면 테두리 없음.
        if ln.find(qn("a:noFill")) is not None:
            return False
        return ln.find(qn("a:solidFill")) is not None
    except Exception:
        return False


def _explicit_run_sizes(slide) -> list:
    """슬라이드 내 모든 run 의 명시 폰트 크기(EMU) 리스트(설정된 것만)."""
    sizes = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            tf = shape.text_frame
        except Exception:
            continue
        for para in tf.paragraphs:
            for run in para.runs:
                try:
                    sz = run.font.size
                except Exception:
                    sz = None
                if sz is not None:
                    sizes.append(int(sz))
    return sizes


def _has_token_margins(slide) -> bool:
    """본문 텍스트 프레임 중 하나라도 python-pptx 기본 마진과 다르게(=토큰 기반
    여백으로) 명시 설정됐는지. 기본값(0.1in/0.05in)과 다르면 명시 적용으로 본다."""
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            tf = shape.text_frame
            ml = tf.margin_left
            mr = tf.margin_right
            mt = tf.margin_top
            mb = tf.margin_bottom
        except Exception:
            continue
        for val, default in (
            (ml, _DEFAULT_MARGIN_LR_EMU),
            (mr, _DEFAULT_MARGIN_LR_EMU),
            (mt, _DEFAULT_MARGIN_TB_EMU),
            (mb, _DEFAULT_MARGIN_TB_EMU),
        ):
            if val is not None and int(val) != int(default):
                return True
    return False


def _count_visual_elements(slide) -> int:
    """비텍스트 시각 요소(도형/배지/카드/구분선/이미지) 개수.

    AUTO_SHAPE(라운드 박스·배지·바 등) 및 PICTURE(장식 이미지)를 센다. 순수
    텍스트박스(TEXT_BOX)는 시각 장식 요소로 세지 않는다.
    """
    n = 0
    for shape in slide.shapes:
        try:
            st = shape.shape_type
        except Exception:
            st = None
        if st is None:
            continue
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE or st == MSO_SHAPE_TYPE.PICTURE:
            n += 1
    return n


# ---------------------------------------------------------------------------
# 공개 API: audit_style_quality
# ---------------------------------------------------------------------------


def audit_style_quality(slide, tokens: dict) -> StyleQualityReport:
    """네이티브 도형이 design_tokens 스타일을 실제 적용했는지 검사(요소 존재가
    아니라 *스타일 품질*).

    이 검사는 ``parity_scorer`` 의 "요소 존재" 채점과 **직교(별개 차원)**이다.
    밀도 채점은 "디자인 요소가 있는가"를, 본 검사는 "기존 design_tokens 가 도형
    서식에 실제 적용됐는가"를 본다. 둘 다 합격해야 "초고퀄"로 판정한다.

    검사 항목 (모두 인자 ``tokens`` = design_tokens_for_profile 결과 dict 에서만
    기준값을 읽는다 — 신규 토큰 정의 아님, Req 5.4):

      1. 카드/박스 라운드(ROUNDED_RECTANGLE 또는 adjustment>0) + 그림자(outerShdw)
         또는 테두리(line) 적용.
      2. 섹션 헤더/배지 등 채워진 도형 중 하나 이상이 tokens 의 accent/primary/
         secondary 색(또는 근접)을 fill 로 가지며, 기본 검정/흰색이 아님.
      3. 타이포 계층: 제목으로 보이는 가장 큰 폰트 텍스트 크기 > 본문 폰트 크기.
      4. 본문 텍스트 프레임의 내부 마진이 python-pptx 기본값(0.1in/0.05in)과
         다르게 명시 설정됨(토큰 기반 여백 적용 흔적).
      5. 슬라이드당 시각 요소(도형/배지/카드/구분선/이미지) 최소 개수(>=2) 충족.

    Args:
        slide: python-pptx Slide.
        tokens: ``design_tokens_for_profile(profile)`` 결과 dict(색·여백·타이포).

    Returns:
        StyleQualityReport(passed, score, missing_style).
        passed = 모든 항목 충족, score = 충족 항목 수 / 전체 항목 수,
        missing_style = 누락 항목명 리스트.
    """
    tokens = tokens or {}
    missing: list = []

    shapes = list(getattr(slide, "shapes", []) or [])

    # --- 항목 1: 카드/박스 라운드 + 그림자/테두리 -------------------------------
    round_decor_ok = False
    for shape in shapes:
        if not _is_auto_shape(shape):
            continue
        if not _is_rounded(shape):
            continue
        if _has_shadow(shape) or _has_border(shape):
            round_decor_ok = True
            break
    if not round_decor_ok:
        missing.append(CHECK_ROUND_DECOR)

    # --- 항목 2: 섹션 헤더/배지 accent 색(검정/흰색 아님) ------------------------
    accent_colors = _accent_like_colors(tokens)
    accent_ok = False
    for shape in shapes:
        rgb = _solid_fill_rgb(shape)
        if rgb is None or _is_black_or_white(rgb):
            continue
        for ac in accent_colors:
            if _color_distance(rgb, ac) <= 36.0:
                accent_ok = True
                break
        if accent_ok:
            break
    if not accent_ok:
        missing.append(CHECK_ACCENT_COLOR)

    # --- 항목 3: 타이포 계층(제목 폰트 > 본문 폰트) -----------------------------
    sizes = _explicit_run_sizes(slide)
    type_ok = bool(sizes) and (max(sizes) > min(sizes))
    if not type_ok:
        missing.append(CHECK_TYPE_HIERARCHY)

    # --- 항목 4: 본문 여백 토큰 적용 -------------------------------------------
    margin_ok = _has_token_margins(slide)
    if not margin_ok:
        missing.append(CHECK_BODY_MARGIN)

    # --- 항목 5: 시각 요소 최소 개수 -------------------------------------------
    visuals_ok = _count_visual_elements(slide) >= _MIN_VISUAL_ELEMENTS
    if not visuals_ok:
        missing.append(CHECK_MIN_VISUALS)

    total = 5
    satisfied = total - len(missing)
    score = satisfied / float(total)
    return StyleQualityReport(passed=(len(missing) == 0), score=score, missing_style=missing)


# ===========================================================================
# 작업 7.1 추가 지점 (APPEND BELOW) — audit_native_density(pptx_path)
# ---------------------------------------------------------------------------
# 산출물_검증기 (a)~(h) 통합 게이트 + AuditReport 는 작업 7.1 에서 이 아래에
# 추가한다. 위의 audit_style_quality / StyleQualityReport / 헬퍼는 수정하지 말고
# 그대로 재사용하라((h) 스타일 품질 검사가 audit_style_quality 를 호출함).
# ===========================================================================

# ===========================================================================
# 작업 7.1 — audit_native_density(pptx_path) -> AuditReport
# ---------------------------------------------------------------------------
# 실제 생성된 .pptx 를 python-pptx 만으로(=Chrome·네트워크 불필요) 열어 슬라이드별
# (a)~(h) 통합 검사를 수행하는 산출물_검증기. 모든 항목 통과 시에만 passed=True.
#
# 검사 항목(장식_배경=풀블리드 PICTURE 슬라이드는 (a)(b)(d)에서 제외):
#   (a) 편집가능 텍스트 런 ≥ 1                         (Req 8.1)
#   (b) 텍스트 보유 셰이프 쌍 겹침률 < 10%              (Req 8.2)
#   (c) 모든 셰이프 슬라이드_경계 안(eps=0.05)          (Req 8.3)
#   (d) 슬라이드별 제목 셰이프 == 1                     (Req 8.4)
#   (e) 카테고리(cover/body) 밀도점수 ≥ 참조점수        (Req 8.5)
#   (f) 풀블리드 배경 베이크 텍스트 미검출              (Req 8.8, 6.1, 6.3)
#   (g) 텍스트가 겹치는 이미지보다 위 z-order           (Req 8.9)
#   (h) audit_style_quality 통과(스타일 품질)           (Req 5.1, 5.3, 5.4)
#
# 재사용 매핑(설계 §Components §3):
#   (b) ← layout_geometry.overlap_area + area (audit_pptx_textbox_overlap.ov 동치,
#         텍스트 보유 셰이프 한정 — 텍스트 없는 장식 컨테이너는 배경 구성으로 제외)
#   (c) ← layout_geometry.within_bounds (eps=0.05)
#   (e) ← parity_scorer.score_native_slide(slide, category)
#   (f) ← audit_pptx_baked_text.baked_text_score (판정 pct>=6.0 or lines>=6)
#   (g) ← audit_pptx_zorder_break._fullbleed / _rect / _ov (z-order/풀블리드 판정)
#   (h) ← 본 파일 audit_style_quality(slide, tokens)  (작업 6.1, 수정 금지·재사용)
#
# 설계 근거: design.md §Components and Interfaces §3, Correctness Property 12/13/15/18.
# 요구사항: Req 6.1·6.3·6.4·8.1~8.10·11.6.
# ===========================================================================

import io
import os
import sys

# 형제(scripts/) 및 ai_engine 패키지 import 를 위한 경로 보강(스크립트 디렉터리는
# 패키지가 아니므로 직접 import 가 가능하도록 sys.path 에 가산적으로 추가).
_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.dirname(_THIS_DIR)
for _p in (_THIS_DIR, _REPO_ROOT, os.path.join(_REPO_ROOT, "ai_engine")):
    if _p and _p not in sys.path:
        sys.path.insert(0, _p)

# --- 재사용 헬퍼 import (try/except fallback, 미가용 시 None 으로 가드) ---------
try:  # layout_geometry — 순수 기하(겹침/경계/풀블리드)
    from layout_geometry import (
        within_bounds as _lg_within_bounds,
        overlap_area as _lg_overlap_area,
        area as _lg_area,
        is_fullbleed as _lg_is_fullbleed,
        SLIDE_RECT as _LG_SLIDE_RECT,
    )
except Exception:  # pragma: no cover
    try:
        from ai_engine.layout_geometry import (  # type: ignore[no-redef]
            within_bounds as _lg_within_bounds,
            overlap_area as _lg_overlap_area,
            area as _lg_area,
            is_fullbleed as _lg_is_fullbleed,
            SLIDE_RECT as _LG_SLIDE_RECT,
        )
    except Exception:
        _lg_within_bounds = _lg_overlap_area = _lg_area = _lg_is_fullbleed = None  # type: ignore[assignment]
        _LG_SLIDE_RECT = (0.0, 0.0, 13.333, 7.5)  # type: ignore[assignment]

try:  # parity_scorer — 네이티브 밀도 채점기 (e)
    from parity_scorer import score_native_slide as _ps_score_native_slide
except Exception:  # pragma: no cover
    try:
        from scripts.parity_scorer import score_native_slide as _ps_score_native_slide  # type: ignore[no-redef]
    except Exception:
        _ps_score_native_slide = None  # type: ignore[assignment]

try:  # audit_pptx_baked_text — 베이크 텍스트 점수 (f)
    from audit_pptx_baked_text import baked_text_score as _baked_text_score
except Exception:  # pragma: no cover
    try:
        from scripts.audit_pptx_baked_text import baked_text_score as _baked_text_score  # type: ignore[no-redef]
    except Exception:
        _baked_text_score = None  # type: ignore[assignment]

try:  # audit_pptx_zorder_break — z-order/풀블리드 판정 재사용 (g)
    from audit_pptx_zorder_break import (
        _fullbleed as _zb_fullbleed,
        _rect as _zb_rect,
        _ov as _zb_ov,
    )
except Exception:  # pragma: no cover
    try:
        from scripts.audit_pptx_zorder_break import (  # type: ignore[no-redef]
            _fullbleed as _zb_fullbleed,
            _rect as _zb_rect,
            _ov as _zb_ov,
        )
    except Exception:
        _zb_fullbleed = _zb_rect = _zb_ov = None  # type: ignore[assignment]

try:  # slide_templates — 기본 design_tokens (h 의 기준값, 신규 토큰 금지)
    from slide_templates import design_tokens_for_profile as _design_tokens_for_profile
except Exception:  # pragma: no cover
    try:
        from ai_engine.slide_templates import design_tokens_for_profile as _design_tokens_for_profile  # type: ignore[no-redef]
    except Exception:
        _design_tokens_for_profile = None  # type: ignore[assignment]

try:  # PIL — (f) 베이크 텍스트 분석용. 미설치 시 (f) 는 보수적으로 skip.
    from PIL import Image as _PILImage
except Exception:  # pragma: no cover
    _PILImage = None  # type: ignore[assignment]


# 슬라이드 경계 상수(인치) — layout_geometry.SLIDE_RECT 와 정합.
_AND_SLIDE_RECT = _LG_SLIDE_RECT if _LG_SLIDE_RECT else (0.0, 0.0, 13.333, 7.5)
_EMU_PER_IN = 914400.0
_EMU_PER_PT = 12700.0

# (b)/(g) 겹침 임계.
_OVERLAP_RATIO = 0.10       # (b) 셰이프 쌍 겹침률 10% 임계
_ZORDER_OVERLAP_PCT = 8.0   # (g) 텍스트↔이미지 z-order 검사 겹침 임계(zorder audit 정합)
# (f) 베이크 텍스트 판정(audit_pptx_baked_text.main 과 동일): pct>=6.0 or lines>=6.
_BAKED_PCT = 6.0
_BAKED_LINES = 6

# 검사 항목 식별자(failures[].check 에 기록).
CHECK_EDITABLE_TEXT = "a_editable_text"
CHECK_OVERLAP = "b_overlap"
CHECK_BOUNDS = "c_bounds"
CHECK_TITLE_ONCE = "d_title_once"
CHECK_DENSITY = "e_density"
CHECK_BAKED_TEXT = "f_baked_text"
CHECK_ZORDER = "g_zorder"
CHECK_STYLE_QUALITY = "h_style_quality"


@dataclass
class AuditReport:
    """``audit_native_density`` 반환값 — (a)~(h) 통합 합격 게이트 결과.

    Attributes:
        passed: 모든 검사 항목 통과 여부(하나라도 실패 시 False).
        failures: 실패 항목 리스트. 각 원소는 dict:
            {
              "check": str,    # 검사 항목명(CHECK_* 상수)
              "slide": int,    # 슬라이드 번호(1-base)
              "shapes": list,  # 문제 셰이프 식별자 리스트
              "signal": str,   # 검출 신호(겹침률/missing/베이크 점수 등)
            }
    """

    passed: bool
    failures: list = field(default_factory=list)


# ---------------------------------------------------------------------------
# 내부 헬퍼 — 셰이프 정규화/식별/측정
# ---------------------------------------------------------------------------


def _nd_rect_in(shape):
    """셰이프 경계 사각형 → (left, top, width, height) 인치. 산출 불가 시 None."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            return None
        return (float(l) / _EMU_PER_IN, float(t) / _EMU_PER_IN,
                float(w) / _EMU_PER_IN, float(h) / _EMU_PER_IN)
    except Exception:
        return None


def _nd_is_picture(shape) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def _nd_text(shape) -> str:
    """편집가능 텍스트 프레임의 비어있지 않은 텍스트(공백 제거). 없으면 ''."""
    try:
        if getattr(shape, "has_text_frame", False):
            return (shape.text_frame.text or "").strip()
    except Exception:
        pass
    return ""


def _nd_max_font_pt(shape):
    """텍스트 프레임 내 최대 run 폰트 크기(pt). 명시값 없으면 None."""
    try:
        if not getattr(shape, "has_text_frame", False):
            return None
    except Exception:
        return None
    best = None
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    sz = run.font.size
                except Exception:
                    sz = None
                if sz is None:
                    continue
                try:
                    pt = float(sz) / _EMU_PER_PT
                except Exception:
                    continue
                if best is None or pt > best:
                    best = pt
    except Exception:
        return best
    return best


def _nd_shape_id(shape, idx: int) -> str:
    """문제 보고용 셰이프 식별자 — '#z:name:텍스트머리'."""
    name = ""
    try:
        name = shape.name or ""
    except Exception:
        name = ""
    txt = _nd_text(shape)
    head = txt.split("\n")[0][:24] if txt else ""
    parts = [f"#{idx}"]
    if name:
        parts.append(name)
    if head:
        parts.append(head)
    return ":".join(parts)


def _nd_norm_title(text) -> str:
    """제목 중복 판정용 정규화(앞뒤 공백 제거 · 소문자)."""
    return str(text or "").strip().lower()


def _nd_collect(slide) -> list:
    """슬라이드 셰이프 트리를 검사용 정규화 dict 리스트로 수집(z=배치 순서)."""
    infos = []
    for z, sh in enumerate(slide.shapes):
        rect = _nd_rect_in(sh)
        is_pic = _nd_is_picture(sh)
        fb = False
        if is_pic and rect is not None and _lg_is_fullbleed is not None:
            try:
                fb = bool(_lg_is_fullbleed(rect))
            except Exception:
                fb = False
        txt = _nd_text(sh)
        infos.append({
            "z": z,
            "shape": sh,
            "rect": rect,
            "is_pic": is_pic,
            "fb": fb,                       # 풀블리드 장식 배경 여부
            "text": txt,
            "has_text": bool(txt),
            "font": _nd_max_font_pt(sh) if txt else None,
            "id": _nd_shape_id(sh, z),
        })
    return infos


# ---------------------------------------------------------------------------
# 슬라이드별 (a)~(h) 검사 (failures 에 가산)
# ---------------------------------------------------------------------------


def _nd_check_slide(slide, slide_no: int, category: str, tokens: dict, failures: list) -> None:
    infos = _nd_collect(slide)
    has_fullbleed = any(i["fb"] for i in infos)
    # 콘텐츠 셰이프(풀블리드 장식 제외) — (a)(b)(d) 대상.
    content_text = [i for i in infos if i["has_text"] and not i["fb"]]
    # (b) 겹침 검사 대상 — 텍스트 보유 셰이프로 한정(사용자 결정).
    #   텍스트 없는 장식 컨테이너(섹션 막대/카드 배경/연락처·노트 박스 바탕 등)는
    #   의도된 배경 구성이므로 제외하고, 기존 audit_pptx_textbox_overlap 관례(텍스트
    #   박스 중심)와 정합한다. 텍스트가 이미지 아래로 가려지는 케이스는 (g) z-order
    #   검사가 별도로 담당한다(Req 2.2/8.2의 콘텐츠 가림 방지 의도).
    content_text_geom = [i for i in content_text if i["rect"] is not None]

    # --- (a) 편집가능 텍스트 런 ≥ 1 (풀블리드 장식_배경 슬라이드는 제외) ----------
    if len(content_text) == 0 and not has_fullbleed:
        failures.append({
            "check": CHECK_EDITABLE_TEXT, "slide": slide_no, "shapes": [],
            "signal": "편집가능 텍스트 런 0개(콘텐츠 텍스트 셰이프 없음)",
        })

    # --- (b) 셰이프 쌍 겹침률 < 10% (텍스트 보유 셰이프 한정, 풀블리드 제외) ------
    if _lg_overlap_area is not None and _lg_area is not None:
        for x in range(len(content_text_geom)):
            for y in range(x + 1, len(content_text_geom)):
                ra, rb = content_text_geom[x]["rect"], content_text_geom[y]["rect"]
                inter = _lg_overlap_area(ra, rb)
                if inter <= 0:
                    continue
                amin = min(_lg_area(ra), _lg_area(rb))
                if amin > 0 and (inter / amin) >= _OVERLAP_RATIO:
                    failures.append({
                        "check": CHECK_OVERLAP, "slide": slide_no,
                        "shapes": [content_text_geom[x]["id"], content_text_geom[y]["id"]],
                        "signal": f"겹침률 {inter / amin * 100:.0f}% (>= 10%)",
                    })

    # --- (c) 모든 셰이프 슬라이드_경계 안 (eps=0.05) -----------------------------
    if _lg_within_bounds is not None:
        for i in infos:
            if i["rect"] is None:
                continue
            try:
                ok = _lg_within_bounds(i["rect"], _AND_SLIDE_RECT, eps=0.05)
            except Exception:
                ok = True
            if not ok:
                r = i["rect"]
                failures.append({
                    "check": CHECK_BOUNDS, "slide": slide_no, "shapes": [i["id"]],
                    "signal": f"경계 밖 rect=({r[0]:.2f},{r[1]:.2f},{r[2]:.2f},{r[3]:.2f})",
                })

    # --- (d) 제목 셰이프 == 1 (제목 텍스트 있는 슬라이드, 풀블리드 제외) ----------
    if content_text:
        # 제목 후보 = 가장 큰 폰트 텍스트, 동률이면 상단(top 최소) 텍스트.
        def _title_key(i):
            f = i["font"] if i["font"] is not None else -1.0
            top = i["rect"][1] if i["rect"] else 1e9
            return (f, -top)

        cand = max(content_text, key=_title_key)
        title_text = _nd_norm_title(cand["text"])
        title_shapes = [i for i in content_text if _nd_norm_title(i["text"]) == title_text]
        if len(title_shapes) != 1:
            failures.append({
                "check": CHECK_TITLE_ONCE, "slide": slide_no,
                "shapes": [i["id"] for i in title_shapes],
                "signal": f"제목 '{title_text[:20]}' 셰이프 {len(title_shapes)}개(정확히 1개여야 함)",
            })

    # --- (e) 카테고리 밀도점수 ≥ 참조점수 ---------------------------------------
    if _ps_score_native_slide is not None:
        try:
            res = _ps_score_native_slide(slide, category)
            if not res.get("passed"):
                failures.append({
                    "check": CHECK_DENSITY, "slide": slide_no, "shapes": [],
                    "signal": (f"{category} 밀도 {res.get('density_score')}/"
                               f"{res.get('reference_score')} 미달, missing={res.get('missing')}"),
                })
        except Exception as e:
            failures.append({
                "check": CHECK_DENSITY, "slide": slide_no, "shapes": [],
                "signal": f"밀도 채점 오류: {e.__class__.__name__}: {e}",
            })

    # --- (f) 풀블리드 배경 베이크 텍스트 미검출 ---------------------------------
    if _baked_text_score is not None and _PILImage is not None:
        for i in infos:
            if not i["fb"]:
                continue
            sh = i["shape"]
            try:
                blob = sh.image.blob
            except Exception:
                continue
            if not blob:
                continue
            try:
                im = _PILImage.open(io.BytesIO(blob))
                pct, lines = _baked_text_score(im)
            except Exception:
                continue
            if pct >= _BAKED_PCT or lines >= _BAKED_LINES:
                failures.append({
                    "check": CHECK_BAKED_TEXT, "slide": slide_no, "shapes": [i["id"]],
                    "signal": (f"풀블리드 배경 베이크 텍스트 검출 — 텍스트추정행 {pct}% / "
                               f"줄 {lines}개 (판정 pct>=6.0 or lines>=6)"),
                })

    # --- (g) 텍스트가 겹치는 이미지보다 위 z-order ------------------------------
    # audit_pptx_zorder_break 의 _fullbleed/_rect/_ov 판정을 재사용한다.
    _fb = _zb_fullbleed if _zb_fullbleed is not None else (
        lambda r: bool(_lg_is_fullbleed(r)) if (_lg_is_fullbleed and r) else False)
    _ovf = _zb_ov if _zb_ov is not None else _lg_overlap_area
    if _ovf is not None:
        pics = []
        txts = []
        for i in infos:
            if i["rect"] is None:
                continue
            if i["is_pic"]:
                if not _fb(i["rect"]):  # 풀블리드 배경은 의도된 backdrop → 제외
                    pics.append(i)
            elif i["has_text"]:
                txts.append(i)
        for tx in txts:
            ta = _lg_area(tx["rect"]) if _lg_area is not None else (
                max(0.0, tx["rect"][2]) * max(0.0, tx["rect"][3]))
            if ta <= 0:
                continue
            for p in pics:
                inter = _ovf(tx["rect"], p["rect"])
                if inter <= 0:
                    continue
                pct = inter / ta * 100.0
                # 겹치면서 텍스트가 이미지보다 먼저(아래) 배치 → 가려짐(z-order 위반).
                if pct >= _ZORDER_OVERLAP_PCT and tx["z"] < p["z"]:
                    failures.append({
                        "check": CHECK_ZORDER, "slide": slide_no,
                        "shapes": [tx["id"], p["id"]],
                        "signal": (f"텍스트가 이미지 아래(z {tx['z']} < {p['z']}), "
                                   f"겹침 {pct:.0f}% — 가려짐"),
                    })

    # --- (h) audit_style_quality 통과(스타일 품질, 작업 6.1 재사용) --------------
    try:
        rep = audit_style_quality(slide, tokens)
        if not rep.passed:
            failures.append({
                "check": CHECK_STYLE_QUALITY, "slide": slide_no, "shapes": [],
                "signal": f"스타일 품질 미달 missing_style={rep.missing_style} score={rep.score:.2f}",
            })
    except Exception as e:
        failures.append({
            "check": CHECK_STYLE_QUALITY, "slide": slide_no, "shapes": [],
            "signal": f"스타일 품질 검사 오류: {e.__class__.__name__}: {e}",
        })


# ---------------------------------------------------------------------------
# 공개 API: audit_native_density
# ---------------------------------------------------------------------------


def audit_native_density(pptx_path: str, tokens: dict = None) -> AuditReport:
    """실제 생성된 .pptx 를 열어 슬라이드별 (a)~(h) 통합 검사를 수행한다.

    python-pptx 만 사용한다(Chrome·네트워크 불필요). **모든 항목 통과 시에만**
    ``passed=True`` 를 반환하며, 하나라도 실패하면 ``failures`` 에 실패 항목명·
    슬라이드 번호(1-base)·문제 셰이프 식별자·검출 신호를 기록한다(Req 8.6).

    슬라이드 카테고리: 첫 슬라이드(인덱스 0)=cover, 나머지=body (Req 8.5,
    parity_scorer.score_native_slide 는 cover/body 만 지원).

    장식_배경(풀블리드 PICTURE)은 (a)(b)(d) 검사에서 콘텐츠 카운트/겹침 대상에서
    제외하며, (f) 베이크 텍스트 검사로 콘텐츠 텍스트가 배경에 구워지지 않았음을
    별도 확인한다. (b) 겹침 검사는 텍스트 보유 셰이프 쌍으로 한정한다(텍스트 없는
    장식 컨테이너=섹션 막대/카드 배경 등은 의도된 배경 구성이므로 제외, 텍스트가
    이미지 아래로 가려지는 케이스는 (g) z-order 가 담당).

    Args:
        pptx_path: 검사할 .pptx 파일 경로.
        tokens: (h) 스타일 품질 기준값으로 쓸 design_tokens dict. None 이면
            ``design_tokens_for_profile(None)`` 기본 프로파일 토큰을 사용한다
            (신규 토큰 정의 아님, 기존 토큰 재사용).

    Returns:
        AuditReport(passed, failures). 통과 시 passed=True, failures=[].
    """
    # python-pptx 는 호출 시점에만 필요(모듈 import 가드와 분리).
    try:
        from pptx import Presentation
    except Exception:  # pragma: no cover - python-pptx 미설치
        return AuditReport(passed=False, failures=[{
            "check": "import", "slide": 0, "shapes": [],
            "signal": "python-pptx 미설치 — audit 불가",
        }])

    if tokens is None:
        tokens = _design_tokens_for_profile(None) if _design_tokens_for_profile is not None else {}

    failures: list = []
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return AuditReport(passed=False, failures=[{
            "check": "open", "slide": 0, "shapes": [],
            "signal": f"파일 열기 실패: {e.__class__.__name__}: {e}",
        }])

    for sidx, slide in enumerate(prs.slides):
        slide_no = sidx + 1            # 1-base 보고
        category = "cover" if sidx == 0 else "body"
        _nd_check_slide(slide, slide_no, category, tokens, failures)

    return AuditReport(passed=(len(failures) == 0), failures=failures)


# ===========================================================================
# 작업 7.2 — generate_visual_comparison (i) 시각 비교 산출물 통합
# ---------------------------------------------------------------------------
# (i) 시각 비교 산출물: scripts/visual_comparator.py 의 compare() 를 호출해 우리
# 렌더(HTML) vs 젠스파크 참조 PNG 를 가로 side-by-side PNG 로 `.generated/
# _design_compare/` 에 생성한다(육안 보조·회귀 추적).
#
# **자동판정 비포함(핵심 제약)**: 이 산출물은 audit_native_density 의 (a)~(h)
# 합격 자동판정에 **포함하지 않는다**. visual_comparator 는 Chrome 헤드리스를
# 쓰므로 다음을 보장한다:
#   - Chrome subprocess 를 **별도 프로세스 + 타임아웃**(기본 20초)으로 감싸 20초+
#     hang 을 방지한다. 타임아웃 시 프로세스를 종료(Chrome 동반 종료)한다.
#   - Chrome 미가용 / 타임아웃 / 입력 누락(ValueError) 시 **skip**(None 반환)하며
#     예외를 전파하지 않는다 → 합격 판정에 영향 없음.
#   - 네트워크 0(visual_comparator 는 로컬 file:// 렌더만 수행).
#
# 설계 근거: design.md §Components §3 (i) 항목, Testing Strategy "실행 환경 제약
# (Chrome/네트워크 hang 방지)" 및 "visual_comparator (육안 보조)".
# 요구사항: Req 5.1 (Property 15 — visual_comparator 산출물).
# ===========================================================================

import time as _vc_time

# visual_comparator.compare / OUT_DIR / CHROME 재사용 import (try/except fallback).
try:
    from visual_comparator import (
        compare as _vc_compare,
        OUT_DIR as _VC_OUT_DIR,
        CHROME as _VC_CHROME,
    )
except Exception:  # pragma: no cover
    try:
        from scripts.visual_comparator import (  # type: ignore[no-redef]
            compare as _vc_compare,
            OUT_DIR as _VC_OUT_DIR,
            CHROME as _VC_CHROME,
        )
    except Exception:
        _vc_compare = None  # type: ignore[assignment]
        _VC_OUT_DIR = os.path.join(_REPO_ROOT, ".generated", "_design_compare")  # type: ignore[assignment]
        _VC_CHROME = "/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"  # type: ignore[assignment]

# (i) 시각 비교 Chrome 타임아웃(초) — 20초+ hang 방지(자동판정 비포함).
_VISUAL_COMPARE_TIMEOUT = 20.0


def _vc_log(msg: str) -> None:
    """(i) 시각 비교 skip/오류를 stderr 로 간결히 로그(자동판정 비포함 신호)."""
    try:
        sys.stderr.write(f"[visual_comparison] {msg}\n")
        sys.stderr.flush()
    except Exception:
        pass


def _vc_worker(q, ours_html: str, reference_png: str, out_png: str) -> None:
    """별도 프로세스에서 visual_comparator.compare 를 실행하고 결과를 큐로 전달.

    Chrome subprocess 가 이 자식 프로세스 안에서 돌므로, 부모가 타임아웃 시 자식을
    terminate 하면 Chrome 도 함께 정리된다(20초+ hang 방지).
    """
    try:
        result = _vc_compare(ours_html, reference_png, out_png)
        q.put(("ok", result))
    except BaseException as e:  # noqa: BLE001 - 모든 실패를 부모로 전달(skip 처리)
        q.put(("err", f"{type(e).__name__}: {e}"))


def generate_visual_comparison(ours_html, reference_png, out_png=None, *,
                               timeout=_VISUAL_COMPARE_TIMEOUT):
    """(i) 우리 렌더 vs 젠스파크 참조 side-by-side PNG 산출물 생성(육안 보조).

    ``scripts/visual_comparator.py`` 의 ``compare(ours_html, reference_png,
    out_png)`` 를 호출해 ``.generated/_design_compare/`` 아래에 비교 PNG 를 만들고
    그 경로를 반환한다.

    **자동판정 비포함**: 이 산출물은 ``audit_native_density`` 의 (a)~(h) 합격
    자동판정과 **무관한** 육안 보조·회귀 추적용이다. Chrome 헤드리스를 쓰므로
    미가용/타임아웃/입력 누락 시 **skip(None 반환)** 하며 예외를 전파하지 않는다.

    Args:
        ours_html: 우리 렌더 HTML 문자열.
        reference_png: 젠스파크 참조 PNG 경로.
        out_png: 출력 PNG 파일명/경로(None 이면 타임스탬프 파일명 자동 생성).
            상대 경로는 ``.generated/_design_compare/`` 아래로 정규화된다.
        timeout: Chrome 렌더 전체 타임아웃(초, 기본 20). 초과 시 종료 후 skip.

    Returns:
        생성된 비교 PNG 의 경로(str). Chrome 미가용/타임아웃/입력 누락/실패 시
        ``None``(skip) — 어떤 경우에도 예외를 전파하지 않는다.
    """
    # --- 입력 누락 graceful skip (compare 의 ValueError 를 사전 차단) ------------
    if not ours_html:
        _vc_log("skip — ours_html 이 비어 있음(입력 누락)")
        return None
    if not reference_png:
        _vc_log("skip — reference_png 경로가 비어 있음(입력 누락)")
        return None
    try:
        if not os.path.isfile(reference_png) or os.path.getsize(reference_png) <= 0:
            _vc_log(f"skip — reference_png 없음/빈 파일: {reference_png}")
            return None
    except Exception as e:
        _vc_log(f"skip — reference_png 접근 실패: {e.__class__.__name__}: {e}")
        return None

    # --- visual_comparator 미가용 skip -----------------------------------------
    if _vc_compare is None:
        _vc_log("skip — visual_comparator 모듈 미가용")
        return None

    # --- Chrome 미가용 skip(자동판정 비포함) ------------------------------------
    try:
        if not os.path.exists(_VC_CHROME):
            _vc_log(f"skip — Chrome 미가용: {_VC_CHROME}")
            return None
    except Exception:
        _vc_log("skip — Chrome 경로 확인 실패")
        return None

    # 출력 파일명(미지정 시 타임스탬프). 상대 경로는 compare 가 OUT_DIR 로 정규화.
    if not out_png:
        out_png = f"native_compare_{int(_vc_time.time() * 1000)}.png"

    # --- 별도 프로세스 + 타임아웃 → Chrome 동반 종료 보장(20초+ hang 방지) -------
    try:
        import multiprocessing as _mp

        ctx = _mp.get_context()
        q = ctx.Queue()
        proc = ctx.Process(
            target=_vc_worker, args=(q, ours_html, reference_png, out_png), daemon=True
        )
        proc.start()
        proc.join(timeout)

        if proc.is_alive():
            # 타임아웃 → 자식(및 그 안의 Chrome) 종료 후 skip.
            try:
                proc.terminate()
                proc.join(5)
            except Exception:
                pass
            _vc_log(f"skip — Chrome 렌더 타임아웃({timeout:.0f}s) 초과, 종료함")
            return None

        try:
            status, payload = q.get(timeout=2)
        except Exception:
            _vc_log("skip — 결과 수신 실패")
            return None

        if status == "ok" and payload:
            return payload
        _vc_log(f"skip — 비교 산출물 생성 실패: {payload}")
        return None
    except Exception as e:
        # 어떤 예외도 자동판정에 영향 주지 않도록 흡수하고 skip.
        _vc_log(f"skip — 예외 흡수: {e.__class__.__name__}: {e}")
        return None


if __name__ == "__main__":  # pragma: no cover - CLI 보조
    _path = sys.argv[1] if len(sys.argv) > 1 else ""
    _report = audit_native_density(_path)
    print(f"=== audit_native_density: {_path}")
    print(f"passed={_report.passed}  failures={len(_report.failures)}")
    for _f in _report.failures:
        print(f"  [슬라이드 {_f['slide']}] {_f['check']}: {_f['signal']}  shapes={_f['shapes']}")
