"""전수 점검: PPTX 슬라이드의 텍스트/이미지 겹침 + 편집 불가(래스터) 텍스트 진단.

사용법:
  ./venv/bin/python scripts/audit_pptx_overlap.py "<pptx 경로>"

분석 항목 (슬라이드별):
  - 모든 shape 의 타입/위치(인치)/크기
  - PICTURE(래스터 이미지): 풀블리드 배경 여부, 화면 점유율
  - 텍스트 프레임(편집 가능 텍스트): 글자 수
  - 텍스트 shape ↔ PICTURE shape 의 사각형 겹침(IoU/overlap%)
  - 슬라이드가 '이미지만'이고 편집 텍스트 0 → 텍스트가 이미지에 구워진(편집 불가) 신호
"""
from __future__ import annotations

import sys
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_PER_IN = 914400.0
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
SLIDE_AREA = SLIDE_W_IN * SLIDE_H_IN


def _in(emu):
    if emu is None:
        return None
    return round(float(emu) / EMU_PER_IN, 3)


def _rect(shape):
    """(left, top, width, height) in inches, or None if missing."""
    try:
        l, t, w, h = shape.left, shape.top, shape.width, shape.height
        if None in (l, t, w, h):
            return None
        return (_in(l), _in(t), _in(w), _in(h))
    except Exception:
        return None


def _overlap_area(r1, r2):
    """두 사각형 (l,t,w,h) 의 겹침 면적(인치^2)."""
    l1, t1, w1, h1 = r1
    l2, t2, w2, h2 = r2
    ix = max(0.0, min(l1 + w1, l2 + w2) - max(l1, l2))
    iy = max(0.0, min(t1 + h1, t2 + h2) - max(t1, t2))
    return ix * iy


def _shape_text(shape):
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        return (shape.text_frame.text or "").strip()
    except Exception:
        return ""


def _is_fullbleed(rect):
    if not rect:
        return False
    l, t, w, h = rect
    # 좌상단 근처에서 시작하고 슬라이드 거의 전체를 덮으면 풀블리드 배경
    return l <= 0.3 and t <= 0.3 and w >= SLIDE_W_IN * 0.92 and h >= SLIDE_H_IN * 0.92


def audit(path):
    prs = Presentation(path)
    print(f"=== 전수 점검: {path}")
    print(f"슬라이드 수: {len(prs.slides)}  (캔버스 {SLIDE_W_IN}x{SLIDE_H_IN} in)\n")

    grand = {
        "slides": len(prs.slides),
        "fullbleed_bg": 0,
        "text_baked_slides": [],     # 이미지만 있고 편집텍스트 0
        "overlap_slides": [],        # 텍스트가 이미지 위에 심하게 겹침
        "pics_total": 0,
        "editable_text_shapes": 0,
    }

    for idx, slide in enumerate(prs.slides, 1):
        pics = []        # (rect, fullbleed, name)
        texts = []       # (rect, chars, name)
        other = 0
        for sh in slide.shapes:
            st = sh.shape_type
            rect = _rect(sh)
            if st == MSO_SHAPE_TYPE.PICTURE:
                fb = _is_fullbleed(rect)
                pics.append((rect, fb, sh.name))
                grand["pics_total"] += 1
                if fb:
                    grand["fullbleed_bg"] += 1
            else:
                txt = _shape_text(sh)
                if txt:
                    texts.append((rect, len(txt), sh.name))
                    grand["editable_text_shapes"] += 1
                else:
                    other += 1

        n_editable_chars = sum(c for _, c, _ in texts)
        fb_pics = [p for p in pics if p[1]]

        print(f"[슬라이드 {idx}] PICTURE={len(pics)} (풀블리드 {len(fb_pics)})  "
              f"편집텍스트 shape={len(texts)} (총 {n_editable_chars}자)  기타 shape={other}")

        # 1) 편집 불가(래스터 텍스트) 신호: 풀블리드 이미지가 있고 편집 텍스트가 없음
        if fb_pics and n_editable_chars == 0:
            grand["text_baked_slides"].append(idx)
            print(f"    ⚠ 편집 불가 신호: 풀블리드 이미지만 존재, 편집 가능한 텍스트 0 "
                  f"→ 텍스트가 이미지에 구워졌을 가능성(편집 불가)")

        # 2) 텍스트 ↔ 이미지 겹침 (풀블리드 배경은 의도된 배경이므로 제외, 부분 이미지만)
        partial_pics = [p for p in pics if not p[1] and p[0]]
        for trect, tchars, tname in texts:
            if not trect:
                continue
            tarea = trect[2] * trect[3]
            if tarea <= 0:
                continue
            for prect, _fb, pname in partial_pics:
                ov = _overlap_area(trect, prect)
                if ov > 0:
                    pct = ov / tarea * 100.0
                    if pct >= 15.0:  # 텍스트 박스의 15%+ 가 부분이미지와 겹침
                        grand["overlap_slides"].append(idx)
                        print(f"    ⚠ 겹침: 텍스트 '{tname}'({tchars}자)의 {pct:.0f}%가 "
                              f"이미지 '{pname}'와 겹침")

        # 3) 풀블리드 위 텍스트 (배경 이미지 위 텍스트 — 가독성 점검 필요)
        if fb_pics and texts:
            print(f"    · 풀블리드 배경 위 편집텍스트 {len(texts)}개 "
                  f"(배경-텍스트 대비/스크림 확인 권장)")

    print("\n=== 요약 ===")
    print(f"총 슬라이드           : {grand['slides']}")
    print(f"총 PICTURE           : {grand['pics_total']}  (풀블리드 배경 {grand['fullbleed_bg']})")
    print(f"편집 가능 텍스트 shape : {grand['editable_text_shapes']}")
    print(f"편집 불가(래스터) 의심 슬라이드 : {sorted(set(grand['text_baked_slides'])) or '없음'}")
    print(f"텍스트·이미지 겹침 슬라이드     : {sorted(set(grand['overlap_slides'])) or '없음'}")

    # 판정 결과를 프로그램적으로 검증할 수 있도록 집계 dict를 반환한다. CLI 사용
    # (아래 __main__)에는 영향이 없다(반환값 미사용). "텍스트·이미지 겹침 슬라이드"는
    # grand["overlap_slides"], "편집 불가(래스터) 의심 슬라이드"는 grand["text_baked_slides"]
    # 로 접근한다(중복 제거 전 원시 리스트 — 개수 판정 시 set()으로 정규화 권장).
    return grand


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("usage: audit_pptx_overlap.py <pptx path>")
        raise SystemExit(2)
    audit(sys.argv[1])
