"""Regression — 직접 PPTX 경로가 Vertex(Nano Banana Pro) 이미지를 네이티브 박스
다이어그램보다 우선 임베드한다.

사용자 보고: 흐름도/구조 슬라이드가 'AI/LNP/Bioinfo' 단어 파편 박스로 떨어지고
Vertex가 전혀 활용되지 않음. 근본 원인: 직접 _tool_generate_pptx 경로에 Vertex 코드
부재 → _classify_section_diagram이 항상 native 박스로 보냄.

수정: native 박스 폴백 *전에* Vertex 키가 활성이면 슬라이드 이미지를 먼저 생성해
임베드한다. Vertex 실패 시에만 네이티브로 폴백.

Correctness properties:
  P1. Vertex 활성 시: 흐름/구조 섹션도 Vertex 이미지(Picture)가 임베드된다.
  P2. Vertex 성공 시 네이티브 박스 다이어그램(_classify_section_diagram)은 호출되지 않는다.
  P3. Vertex 실패(error) 시: 네이티브 경로로 안전 폴백한다(에러 없이 생성 완료).
  P4. AE_PREFER_VERTEX_IMAGE=0 이면 Vertex를 건너뛴다.

실행: pytest scripts/test_pptx_vertex_first.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
vmod = pytest.importorskip("ai_engine.vertex_image_module")
pptx = pytest.importorskip("pptx")

# 1x1 PNG
_PNG_B64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAACklEQVR4nGNgAAAAAgAB"
    "c3IAAAAASUVORK5CYII=")


class _FakeVertex:
    enabled = True
    last_prompts = []

    def __init__(self, fail=False):
        self._fail = fail

    async def generate(self, **kwargs):
        _FakeVertex.last_prompts.append(kwargs.get("prompt", ""))
        if self._fail:
            return {"error": "quota", "detail": "test"}
        return {"images": [_PNG_B64], "model": "gemini-3-pro-image-preview"}


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_vertex_image_preferred_over_native(tmp_path, monkeypatch):
    """P1+P2 — Vertex 활성 시 흐름/구조 섹션이 Vertex 이미지로 임베드되고 네이티브 미호출."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "1")
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "0")  # image 모드: Vertex 사전루프 활성
    # HTML 브리지는 비활성으로(네이티브/Vertex 분기만 검증)
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)
    monkeypatch.setattr(vmod, "get_vertex_image_client", lambda *a, **k: _FakeVertex())

    # 네이티브 분류기가 호출되면 실패 — Vertex가 우선되어야 한다
    def _boom(*a, **k):
        raise AssertionError("native _classify_section_diagram should NOT run when Vertex succeeds")
    monkeypatch.setattr(server, "_classify_section_diagram", _boom)

    slides = [
        {"title": "EC2 비용 모니터링 에이전트 흐름", "bullets": ["pynvml GPU 감지 -> psutil 스캔 -> PID 추적"]},
        {"title": "MOGAM 멀티 에이전트 시스템", "bullets": ["Team Lead 라우터 -> AI 분류 -> LNP/Bioinfo"]},
    ]
    out = asyncio.run(server._tool_generate_pptx({"title": "프로젝트 흐름도 분석", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        # cover + 2 content = 3
        assert len(prs.slides._sldIdLst) == 3
        # 콘텐츠 슬라이드(2,3)에 Vertex 이미지(Picture)가 있어야 한다
        content_with_pic = 0
        for sl in list(prs.slides)[1:]:
            if any(sh.shape_type == MSO_SHAPE_TYPE.PICTURE for sh in sl.shapes):
                content_with_pic += 1
        assert content_with_pic >= 2, f"Vertex 이미지가 임베드된 콘텐츠 슬라이드 수={content_with_pic}"
    finally:
        _cleanup(res)


def test_falls_back_to_native_when_vertex_fails(tmp_path, monkeypatch):
    """P3 — Vertex 실패 시 네이티브 경로로 폴백(에러 없이 생성)."""
    from pptx import Presentation
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "1")
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "0")  # image 모드: Vertex 사전루프 활성
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)
    monkeypatch.setattr(vmod, "get_vertex_image_client", lambda *a, **k: _FakeVertex(fail=True))

    called = {"native": 0}
    real_classify = server._classify_section_diagram

    def _spy(h, b, t=""):
        called["native"] += 1
        return real_classify(h, b, t)
    monkeypatch.setattr(server, "_classify_section_diagram", _spy)

    slides = [{"title": "구조 분석", "bullets": ["A -> B -> C"]}]
    out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        # Vertex 실패 → 네이티브 분류기가 호출되어 폴백 동작
        assert called["native"] >= 1, "Vertex 실패 시 네이티브 폴백이 동작해야 함"
        prs = Presentation(res["absPath"])
        assert len(prs.slides._sldIdLst) == 2
    finally:
        _cleanup(res)


def test_opt_out_skips_vertex(tmp_path, monkeypatch):
    """P4 — AE_PREFER_VERTEX_IMAGE=0 이면 Vertex를 건너뛰고 네이티브로 간다."""
    from pptx import Presentation
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "0")
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "0")
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)

    def _boom(*a, **k):
        raise AssertionError("Vertex client must not be fetched when opted out")
    monkeypatch.setattr(vmod, "get_vertex_image_client", _boom)

    slides = [{"title": "개요", "bullets": ["항목 1", "항목 2"]}]
    out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        prs = Presentation(res["absPath"])
        assert len(prs.slides._sldIdLst) == 2
    finally:
        _cleanup(res)


def test_visual_intent_keywords_expanded():
    """흐름도/프로세스/아키텍처 등이 visual_intent를 유발한다(고품질 tier 활성화)."""
    for kw in ["프로젝트 흐름도 분석", "시스템 아키텍처", "데이터 파이프라인 구성",
               "process flow", "발표 자료 개요"]:
        assert server._detect_visual_intent(kw), f"visual_intent 미감지: {kw}"


def test_vertex_prompt_is_content_faithful(tmp_path, monkeypatch):
    """개선 — Vertex 프롬프트가 실제 라벨 + 다이어그램 유형 지시를 담는다(젠스파크 충실형).

    장식 일러스트가 아니라 '실제 라벨이 들어간 다이어그램'을 그리도록 프롬프트를
    강화했다. 흐름(→/단계) 섹션이면 flowchart 지시 + 정확한 한글 라벨이 포함돼야 한다."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    monkeypatch.setenv("AE_PREFER_VERTEX_IMAGE", "1")
    monkeypatch.setenv("AE_PREFER_EDITABLE_DIAGRAM", "0")  # image 모드: Vertex 사전루프 활성
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: None)
    _FakeVertex.last_prompts = []
    monkeypatch.setattr(vmod, "get_vertex_image_client", lambda *a, **k: _FakeVertex())
    monkeypatch.setattr(server, "_classify_section_diagram",
                        lambda *a, **k: (_ for _ in ()).throw(AssertionError("native must not run")))

    slides = [{"title": "데이터 처리 흐름",
               "bullets": ["수집 단계", "정제 단계", "적재 단계"]}]
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "프로젝트 흐름도 분석", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, res
        # 콘텐츠 슬라이드(2번째)에 대한 프롬프트 확인
        joined = "\n".join(_FakeVertex.last_prompts)
        assert "flowchart" in joined, "흐름 섹션인데 flowchart 지시가 없음"
        # 실제 한글 라벨이 프롬프트에 verbatim 포함
        assert "수집 단계" in joined and "적재 단계" in joined, "실제 라벨이 프롬프트에 없음"
        assert "NO watermark" in joined and "16:9" in joined
    finally:
        _cleanup(res)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
