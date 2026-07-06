"""Regression — 슬라이드 레이아웃 품질(이미지 fit·다이어그램 전체폭/정렬).

배경: 생성 슬라이드의 이미지가 작게/넘쳐서 박히고, 다이어그램이 우측 절반에
눌려 텍스트·행 정렬이 어긋났다. 수정:
  - 이미지: 비율 보존 fit + 영역 중앙정렬(슬라이드 경계 7.5in 안)
  - 다이어그램: 제목 아래 전체 폭(좌 0.6in 시작) 사용 + 중복 bullets 제거
  - flow 5단계 이상: 세로 흐름(넓은 박스)

대상: ai_engine.native_diagram_pptx.build_native_diagram + _tool_generate_pptx
실행: pytest scripts/test_pptx_layout_quality.py -q
"""
from __future__ import annotations

import asyncio
import os
import sys
import tempfile

import pytest

sys.path.insert(0, os.path.abspath(os.path.join(os.path.dirname(__file__), "..")))

pptx = pytest.importorskip("pptx")
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

ndp = pytest.importorskip("ai_engine.native_diagram_pptx")


def _new_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    return prs, slide


def test_flow_vertical_for_many_steps_full_width():
    prs, slide = _new_slide()
    content = " -> ".join([f"단계{i}" for i in range(1, 7)])  # 6단계 → 세로
    ok = ndp.build_native_diagram(slide, "flow", content, region=(0.6, 1.7, 12.1, 5.2))
    assert ok is True
    shapes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(shapes) >= 6, "6단계 흐름 박스가 모두 그려져야 함"
    # 모든 박스가 슬라이드 경계(13.333 x 7.5in) 안에 있어야 함
    for s in shapes:
        assert s.left >= 0 and s.top >= 0
        assert s.left + s.width <= Inches(13.333) + Inches(0.05)
        assert s.top + s.height <= Inches(7.5) + Inches(0.05)


def test_flow_boxes_left_aligned_vertical():
    # 세로 흐름의 박스들은 같은 x(좌표)로 정렬되어야 한다(행 정렬).
    prs, slide = _new_slide()
    content = " -> ".join([f"S{i}" for i in range(1, 6)])
    ndp.build_native_diagram(slide, "flow", content, region=(0.6, 1.7, 12.1, 5.2))
    lefts = [s.left for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert lefts, "박스 없음"
    # 모든 좌측 좌표가 동일(세로 정렬)
    assert max(lefts) - min(lefts) <= Inches(0.02), "세로 흐름 박스 x정렬 어긋남"


def test_image_fits_within_slide_bounds():
    # 정사각 이미지가 슬라이드 하단(7.5in)을 넘치지 않아야 한다.
    server = pytest.importorskip("ai_engine.server")
    from PIL import Image as _PIL
    with tempfile.TemporaryDirectory() as td:
        os.environ["AE_GENERATED_ROOT"] = td
        # 정사각 PNG 생성 후 imageFile로 직접 주입(게이트웨이 호출 회피)
        img_path = os.path.join(td, "sq.png")
        _PIL.new("RGB", (1024, 1024), (80, 120, 200)).save(img_path)
        out = asyncio.run(server._tool_generate_pptx(
            {"title": "img", "slides": [
                {"title": "사진 슬라이드", "bullets": [], "imageFile": img_path},
            ]}, td))
        import re
        m = re.search(r"([^\s'\"]+\.pptx)", out)
        assert m, out[:300]
        path = m.group(1)
        if not os.path.isabs(path):
            path = os.path.join(td, path)
        prs = Presentation(path)
        for slide in prs.slides:
            for s in slide.shapes:
                if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    assert s.top + s.height <= Inches(7.5) + Inches(0.05), "이미지가 슬라이드 하단을 넘침"
                    assert s.left + s.width <= Inches(13.333) + Inches(0.05)
                    # 너무 작지 않아야 함(최소 3인치 폭)
                    assert s.width >= Inches(3.0), "이미지가 너무 작게 배치됨"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
