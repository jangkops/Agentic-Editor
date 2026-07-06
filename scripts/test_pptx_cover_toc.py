"""Regression — 슬라이드 덱의 기본인 표지 + 목차(TOC)가 자동 생성된다.

사용자 보고: "표지/목차가 생성 요청해도 없음. 슬라이드 작업의 기본인데."
또한 HTML 통짜 슬라이드(배경 PNG 위 텍스트 겹침)는 기본 OFF여야 한다.

Correctness properties:
  P1. 콘텐츠 3장 이상이면 '목차' 슬라이드가 표지 다음(2번째)에 생성된다.
  P2. 목차에 각 섹션 제목이 편집 가능한 텍스트로 들어간다(통짜 이미지 아님).
  P3. 표지(1번째)에 덱 제목이 편집 가능 텍스트로 들어간다.
  P4. 기본적으로 슬라이드에 통짜 PNG(Picture)가 없다(편집 가능 우선).

실행: pytest scripts/test_pptx_cover_toc.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def test_cover_and_toc_generated(tmp_path):
    """P1+P2+P3 — 표지 + 목차 자동 생성, 섹션 제목 편집 가능 텍스트로 수록."""
    from pptx import Presentation

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    # 통짜 PNG/네이티브 다이어그램 모두 끄고 표지/목차/텍스트 구조만 검증
    os.environ["AE_DISABLE_NATIVE_DIAGRAM"] = "1"
    os.environ["AE_PPTX_TOC"] = "1"  # 이 테스트는 TOC 동작을 검증
    slides = [
        {"title": "프로젝트 개요", "bullets": ["배경"]},
        {"title": "구조 분석", "bullets": ["디렉토리"]},
        {"title": "핵심 기술", "bullets": ["스택"]},
        {"title": "결론", "bullets": ["요약"]},
    ]
    try:
        out = asyncio.run(server._tool_generate_pptx(
            {"title": "프로젝트 분석 보고서", "slides": slides}, ""))
        res = json.loads(out)
        try:
            assert "error" not in res, res
            prs = Presentation(res["absPath"])
            slist = list(prs.slides)
            # cover + toc + 4 content = 6
            assert len(slist) >= 6, f"슬라이드 수 {len(slist)} (표지+목차+4 기대)"

            def _text(sl):
                return "\n".join(sh.text_frame.text for sh in sl.shapes if sh.has_text_frame)

            # P3 — 표지에 덱 제목
            assert "프로젝트 분석 보고서" in _text(slist[0]), "표지 제목 없음"
            # P1+P2 — 2번째 슬라이드가 목차, 섹션 제목 수록
            toc_text = _text(slist[1])
            assert "목차" in toc_text, "목차 슬라이드 없음"
            for t in ("프로젝트 개요", "구조 분석", "핵심 기술", "결론"):
                assert t in toc_text, f"목차에 섹션 누락: {t}"
        finally:
            _cleanup(res)
    finally:
        os.environ.pop("AE_DISABLE_NATIVE_DIAGRAM", None)
        os.environ["AE_PPTX_TOC"] = "0"  # conftest 기본값으로 복원(누수 방지)


def test_html_slides_off_by_default(tmp_path, monkeypatch):
    """P4 — HTML 통짜 슬라이드가 기본 비활성(브리지 가용해도 PNG 배경 미생성)."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)
    os.environ["AE_DISABLE_NATIVE_DIAGRAM"] = "1"
    # 브리지 가용으로 위장 + HTML 렌더가 호출되면 실패
    monkeypatch.setattr(server, "_call_bridge", lambda *a, **k: {"remote": False})

    called = {"n": 0}

    async def _fake_html(*a, **k):
        called["n"] += 1
        return ".generated/should-not-be-used.png"
    monkeypatch.setattr(server, "_generate_html_slide_for_section", _fake_html)

    slides = [{"title": "개요", "bullets": ["항목"]}, {"title": "상세", "bullets": ["설명"]}]
    try:
        out = asyncio.run(server._tool_generate_pptx({"title": "덱", "slides": slides}, ""))
        res = json.loads(out)
        try:
            assert "error" not in res, res
            assert called["n"] == 0, "HTML 슬라이드가 기본 비활성이어야 하는데 호출됨"
            prs = Presentation(res["absPath"])
            pics = sum(1 for sl in prs.slides for sh in sl.shapes
                       if sh.shape_type == MSO_SHAPE_TYPE.PICTURE)
            assert pics == 0, f"통짜 PNG가 {pics}개 — 기본은 편집 가능이어야 함"
        finally:
            _cleanup(res)
    finally:
        os.environ.pop("AE_DISABLE_NATIVE_DIAGRAM", None)


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
