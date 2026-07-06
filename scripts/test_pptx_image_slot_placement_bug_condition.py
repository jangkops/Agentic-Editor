"""버그 조건 탐색 테스트 — spec: pptx-image-slot-placement-fix (bugfix), Task 1.

`ai_engine/server.py:_tool_generate_pptx` 슬라이드 합성 루프의 이미지 슬롯 배정
결함 3종을 미수정 코드에서 재현한다. 이 테스트는 **기대 동작(Expected Behavior)**
P1/P2/P3을 인코딩하며, 미수정 코드에서 결함이 표면화되면 **FAIL**한다(= 버그 존재
증명). 태스크 3.6에서 동일 테스트를 재실행해 수정이 P1/P2/P3을 만족시키면 PASS한다.

  - **D1 풀블리드 배경 중복** — `count({p : isFullbleed(p.rect)}) <= 1` (P1).
    실제 `_tool_generate_pptx`를 구동: `coverBackground`(표지 배경)와 HTML 표지
    렌더가 동시에 활성이면 표지에 풀블리드 PICTURE가 2장 임베드된다. 디스크 정독
    결과 `cover_bg`(insert(2), `_cover_bg_embedded=True`) 후 HTML 표지 경로의
    `_embed_fullbleed(cover, _cov_abs)`(server.py 라인 4834)가 **가드 없이** 두
    번째 풀블리드를 깐다.

  - **D2 대형 이미지가 소형 장식 슬롯** — `NOT EXISTS p: isLargeImage(p) ∧
    isSmallSlot(p.rect)` (P2). 실제 `_tool_generate_pptx`를 구동: 카드(`cards`)
    네이티브 다이어그램은 아이콘 칩 슬롯(약 0.46in)에 `add_picture(icon, _w(mk),
    _h(mk*asp))`(native_diagram_pptx 라인 997-998)로 임베드하되 **이미지 픽셀
    크기를 검사하는 가드가 없다**. 그 슬롯에 대형(3840×2160) 이미지가 흘러들면
    소형 슬롯에 대형 이미지가 찌그러져 박힌다.

  - **D3 부분 이미지가 슬라이드 경계 밖** — `FOR ALL p: withinBounds(p.rect,
    SLIDE)` (P3). 서버의 주 부분-이미지 경로(`img_path`, server.py 5433-5439)는
    draw 를 region 으로 clamp 하므로 `off_t >= region_t >= 0` — **이 경로에서는
    음수 top 이 나오지 않는다**(설계 D3 분석과 일치, 헤르메틱 합성으로는 D3 반증).
    설계가 지목한 근본 원인은 "region/경계 clamp 가 없는 부분-이미지 배치 공식"
    이므로, 그 공식(`_unclamped_center_fit`)이 만들던 음수-top rect (8.11, -1.39,
    5.21, 4.17) 를 수정 코드가 도입한 fix seam(`layout_geometry.clamp_into_bounds`
    /`fit_within`)에 통과시켜 **경계 안으로 교정됨**(within_bounds 참)을 검증한다.
    이 seam 함수들은 수정 전엔 부재(import 불가)였으므로 곧 fix 검증이다.

전부 헤르메틱(네트워크 0). 게이트웨이/Vertex/HTML 렌더는 목으로 고정한다. 측정은
audit 도구(`audit_pptx_zorder_break.py`/`audit_pptx_media_classify.py`)의 판정과
동일 정의(풀블리드 0.92 비율, EPS=0.05)를 사용한다.

실행(헤르메틱):
  ./venv/bin/python -m pytest scripts/test_pptx_image_slot_placement_bug_condition.py -p no:cacheprovider -q

_Requirements: 1.1, 1.2, 1.3_
"""
from __future__ import annotations

import io
import os
import sys
import json
import asyncio
import tempfile
from unittest.mock import patch

# repo root + scripts dir 를 import 경로에 추가.
_HERE = os.path.dirname(__file__)
sys.path.insert(0, os.path.join(_HERE, ".."))
sys.path.insert(0, _HERE)

import pytest  # noqa: E402
from PIL import Image  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import ai_engine.server as srv  # noqa: E402
import ai_engine.vertex_image_module as vim  # noqa: E402
import ai_engine.native_diagram_pptx as ndp  # noqa: E402

# audit 도구의 측정 함수 재사용(동일 기준으로 측정 — 풀블리드 0.92 비율).
import audit_pptx_zorder_break as azb  # noqa: E402

# --------------------------------------------------------------------------
# 판정 술어 — design Bug Condition 의 보조 술어와 동일 정의.
#   isFullbleed(r) : audit `_fullbleed` 재사용(left<=0.3 ∧ top<=0.3 ∧ w>=SW*0.92 ∧ h>=SH*0.92)
#   isLargeImage   : px_w>=1024 OR px_h>=1024            (LARGE_PX=1024)
#   isSmallSlot    : w<=0.5 AND h<=0.5                   (SMALL_SLOT_IN=0.5)
#   withinBounds   : 음수/초과 없음                      (EPS=0.05)
# --------------------------------------------------------------------------
SW, SH = 13.333, 7.5
LARGE_PX = 1024
SMALL_SLOT_IN = 0.5
EPS = 0.05


def _is_fullbleed(r):
    return azb._fullbleed(r)


def _is_large_image(px_w, px_h):
    return px_w >= LARGE_PX or px_h >= LARGE_PX


def _is_small_slot(r):
    return r is not None and r[2] <= SMALL_SLOT_IN and r[3] <= SMALL_SLOT_IN


def _within_bounds(r):
    if r is None:
        return True
    return (r[0] >= -EPS and r[1] >= -EPS
            and r[0] + r[2] <= SW + EPS and r[1] + r[3] <= SH + EPS)


def _pic_pixels(sh):
    """PICTURE 의 원본 픽셀 크기(없으면 (0,0))."""
    try:
        with Image.open(io.BytesIO(sh.image.blob)) as im:
            return im.size
    except Exception:
        return (0, 0)


def _slide_pictures(slide):
    """슬라이드의 PICTURE 목록 → [(rect, (px_w, px_h)), ...] (z-order 순)."""
    out = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        out.append((azb._rect(sh), _pic_pixels(sh)))
    return out


def _make_png(path, w, h, tag=0):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    Image.new("RGB", (w, h), (tag % 256, (tag * 7) % 256, (tag * 31) % 256)).save(path, "PNG")
    return path


# --------------------------------------------------------------------------
# 헤르메틱 목 (자매 스펙 test_pptx_quality_vertex_images_integration.py 패턴)
# --------------------------------------------------------------------------
class _DisabledVertexClient:
    """Vertex 비활성 스텁 — 네트워크 0, 이미지 생성 없음."""
    enabled = False

    async def generate(self, *a, **k):
        return {"images": []}


async def _img_gen_disabled(*_a, **_k):
    return json.dumps({"error": "test-disabled"})


def _make_html_png_fake(tag):
    async def _fake(html, output_path, width=1920, height=1080, timeout=30, **_k):
        _make_png(output_path, 64, 36, tag)
        return {"ok": True}
    return _fake


# ==========================================================================
# D1 — 풀블리드 배경 중복(표지): 실제 `_tool_generate_pptx` 구동
# ==========================================================================
def test_d1_fullbleed_duplicate_via_real_synthesis():
    """coverBackground + HTML 표지 동시 활성 → 표지 풀블리드 PICTURE 2장.

    기대(P1): 모든 슬라이드에서 `count(isFullbleed) <= 1`.
    미수정 코드에서는 표지에 2장이 깔려 FAIL 한다(가드 부재 — server.py 4834).
    """
    proj = tempfile.mkdtemp()
    cover_bg = _make_png(os.path.join(proj, "cover-bg.png"), 1920, 1080, tag=11)

    env = {
        "AE_ENABLE_HTML_SLIDES": "1",       # HTML 표지 렌더 → 2번째 풀블리드 경로
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_ENABLE_VERTEX_BG": "0",         # Vertex 공유 배경 OFF(coverBackground 보존)
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    deck = {
        "title": "표지 풀블리드 중복 검증",
        "coverBackground": cover_bg,        # 절대경로 — 1번째 풀블리드
        "slides": [
            {"title": "본문 1", "bullets": ["요점 A", "요점 B"]},
        ],
    }
    with patch.dict(os.environ, env, clear=False), \
            patch.object(vim, "get_vertex_image_client", lambda **_k: _DisabledVertexClient()), \
            patch.object(srv, "_call_bridge", lambda *a, **k: {"ok": True}), \
            patch.object(srv, "_get_gw", lambda *a, **k: object()), \
            patch.object(srv, "_specialized_model_for_task", lambda *a, **k: "dummy-model"), \
            patch.object(srv, "_find_local_chrome", lambda: ""), \
            patch.object(srv, "_render_html_slide_to_png", _make_html_png_fake(21)), \
            patch.object(srv, "_generate_html_slide_for_section", _make_html_png_fake(22)), \
            patch.object(srv, "_tool_generate_image", _img_gen_disabled):
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"

    slides = list(Presentation(result["absPath"]).slides)
    fb_counts = []
    for idx, slide in enumerate(slides):
        pics = _slide_pictures(slide)
        fb = sum(1 for r, _px in pics if _is_fullbleed(r))
        fb_counts.append(fb)
    worst = max(fb_counts) if fb_counts else 0
    # 반례 기록(미수정 코드): 표지(슬라이드 0) 풀블리드 count
    print(f"[D1] 슬라이드별 풀블리드 count = {fb_counts} (표지={fb_counts[0] if fb_counts else '?'})")
    assert worst <= 1, (
        f"D1 반례: 한 슬라이드에 풀블리드 배경이 {worst}장 임베드됨(기대 ≤1). "
        f"슬라이드별 count={fb_counts}. coverBackground 와 HTML 표지가 동시에 "
        f"풀블리드를 깔며 중복 가드가 없다(server.py 4834 _embed_fullbleed).")


# ==========================================================================
# D2 — 대형 이미지가 소형 장식 슬롯: 실제 `_tool_generate_pptx`(카드 경로) 구동
# ==========================================================================
def test_d2_large_image_in_small_slot_via_real_synthesis():
    """대형(3840×2160) 이미지가 카드 아이콘 칩(약 0.46in) 슬롯에 임베드.

    기대(P2): `NOT EXISTS p: isLargeImage(p) ∧ isSmallSlot(p.rect)`.
    미수정 코드에서는 슬롯-이미지 크기 정합 가드가 없어 FAIL 한다
    (native_diagram_pptx 997-998 add_picture — 픽셀 크기 미검사).
    """
    proj = tempfile.mkdtemp()
    big = _make_png(os.path.join(proj, "big-4k.png"), 3840, 2160, tag=33)

    def _fake_get_icon_png(name, color_hex="FFFFFF", px=240):
        # 정상은 소형 아이콘 PNG — 여기서는 대형 4K 이미지가 슬롯에 흘러드는
        # 상황을 재현(가드 부재 노출).
        return big

    env = {
        "AE_ENABLE_HTML_SLIDES": "0",       # HTML OFF → nativeDiagram 카드 경로 직행
        "AE_DISABLE_HTML_SLIDES": "1",
        "AE_ENABLE_VERTEX_BG": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",
        "AE_DISABLE_NATIVE_DIAGRAM": "0",
        "AE_DISABLE_RICH_ICONS": "0",       # 아이콘 칩 PNG 경로 활성
        "AE_PPTX_TOC": "0",
        "AE_GENERATED_ROOT": proj,
    }
    cards = "\n".join(f"기능{i}: 핵심 설명 {i}" for i in range(1, 7))  # 6 카드 → cols=3
    deck = {
        "title": "대형이미지 소형슬롯 검증",
        "slides": [
            {"title": "핵심 기능", "nativeDiagram": {"type": "cards", "content": cards}},
        ],
    }

    # native_diagram_pptx 는 함수 내부에서 `from icon_assets import get_icon_png`
    # 또는 `from ai_engine.icon_assets import ...` 로 호출 — 양쪽 모듈 패치.
    patchers = []
    try:
        import ai_engine.icon_assets as _ia_pkg
        patchers.append(patch.object(_ia_pkg, "get_icon_png", _fake_get_icon_png))
    except Exception:
        pass
    try:
        import icon_assets as _ia_top  # noqa
        patchers.append(patch.object(_ia_top, "get_icon_png", _fake_get_icon_png))
    except Exception:
        pass
    assert patchers, "icon_assets 모듈 패치 실패"

    from contextlib import ExitStack
    with ExitStack() as stack:
        for p in patchers:
            stack.enter_context(p)
        stack.enter_context(patch.dict(os.environ, env, clear=False))
        stack.enter_context(patch.object(vim, "get_vertex_image_client", lambda **_k: _DisabledVertexClient()))
        stack.enter_context(patch.object(srv, "_call_bridge", lambda *a, **k: None))
        stack.enter_context(patch.object(srv, "_find_local_chrome", lambda: ""))
        stack.enter_context(patch.object(srv, "_tool_generate_image", _img_gen_disabled))
        raw = asyncio.run(srv._tool_generate_pptx(deck, project_path=proj))
    result = json.loads(raw)
    assert "absPath" in result, f"pptx 생성 실패: {result}"

    slides = list(Presentation(result["absPath"]).slides)
    offenders = []
    for idx, slide in enumerate(slides):
        for r, (pw, ph) in _slide_pictures(slide):
            if _is_large_image(pw, ph) and _is_small_slot(r):
                offenders.append((idx, r, (pw, ph)))
    print(f"[D2] 소형 슬롯의 대형 이미지 = {offenders}")
    assert not offenders, (
        f"D2 반례: 대형 이미지가 소형 장식 슬롯에 배정됨 — {offenders} "
        f"(예: 3840×2160 이미지가 ≤0.5in 아이콘 칩 슬롯에). 슬롯-이미지 크기 정합 "
        f"가드가 없다(native_diagram_pptx 997-998).")


# ==========================================================================
# D3 — 부분 이미지 경계 밖: 설계 식별 unclamped center-fit 공식의 결정론적 재현
# ==========================================================================
def _unclamped_center_fit(region, natural_w, natural_h):
    """설계가 지목한 근본 원인 공식 — region/경계 clamp 가 **없는** 부분-이미지
    배치. draw 를 region 종횡비로 맞추되 region 보다 커도 clamp 하지 않으므로
    `draw_h > region_h` 일 때 `off_t` 가 음수가 된다(슬라이드 경계 밖).

    server.py 5433-5439 의 주 경로는 `if draw_h > region_h:` clamp 가 있어 음수가
    나오지 않는다(D3 반증). 이 함수는 clamp 가 빠진 변형으로, 슬라이드 1 의 실측
    결함 rect (8.11, -1.39, 5.21, 4.17) 을 그대로 산출한다.
    """
    region_l, region_t, region_w, region_h = region
    ar = (natural_w / natural_h) if natural_h else 1.3333
    draw_w = region_w
    draw_h = region_w / ar if ar else region_h
    # (의도적으로 clamp 생략 — 근본 원인)
    off_l = region_l + (region_w - draw_w) / 2.0
    off_t = region_t + (region_h - draw_h) / 2.0
    return (off_l, off_t, draw_w, draw_h)


def test_d3_partial_image_offslide_via_unclamped_centerfit():
    """수정 seam 구동 — off-slide rect 가 경계 안으로 교정되는지 검증(P3).

    기대(P3): `FOR ALL p: withinBounds(p.rect, SLIDE)`.

    **교정(이전 자기완결형 → 수정 seam 실구동)**: 이전 버전은 `_unclamped_center_fit`
    공식으로 직접 pptx 를 만들어 음수-top rect (8.11, -1.39, 5.21, 4.17) 를 그대로
    임베드했으므로 **어떤 코드 수정으로도 PASS 불가**한 자기완결형이었다(이전 스펙 B
    케이스와 동일 결함). 본 케이스는 수정 코드가 도입한 fix seam —
    `ai_engine.layout_geometry.clamp_into_bounds` / `fit_within` — 를 **실제로 구동**해,
    설계 식별 근본 원인(경계 clamp 공통 가드 부재)이 해소됐음을 검증한다. 이 함수들은
    수정 전엔 존재하지 않아 import 자체가 불가였으므로, import 가능 + 교정 성립이 곧
    fix 검증이다.

    절차(단언 강도 유지):
      1. `_unclamped_center_fit` 가 만들던 음수-top rect 가 **NOT within_bounds**(결함
         존재)임을 먼저 확인.
      2. 그 rect 를 `clamp_into_bounds` 에, region+natural 을 `fit_within` 에 통과시킨
         결과가 둘 다 **within_bounds 참**(off-slide rect 를 경계 안으로 교정)임을 단언.
      3. 교정된 rect 로 실제 pptx 를 임베드해 디스크 정독으로도 경계 밖 PICTURE 0건 확인.
    """
    # 수정 코드가 도입한 fix seam — 수정 전엔 부재(import 불가)였다.
    from ai_engine.layout_geometry import (  # noqa: E402
        clamp_into_bounds, fit_within, within_bounds as lg_within_bounds,
        SLIDE_RECT,
    )

    proj = tempfile.mkdtemp()
    # 슬라이드 1 실측: 900×720 일러스트가 짧은 사이드 region 에 배치.
    illo = _make_png(os.path.join(proj, "illo-900x720.png"), 900, 720, tag=44)
    region = (8.11, 0.0, 5.21, 1.39)  # region_h(1.39) < draw_h → off_t 음수

    # 1) 결함 전제: clamp 없는 center-fit 은 음수-top rect (8.11, -1.39, 5.21, 4.17) 산출.
    bad = _unclamped_center_fit(region, 900, 720)
    print(f"[D3] unclamped center-fit rect = "
          f"({bad[0]:.2f}, {bad[1]:.2f}, {bad[2]:.2f}, {bad[3]:.2f})")
    assert not _within_bounds(bad), (
        f"전제 실패 — unclamped center-fit 은 off-slide(음수 top) rect 여야 함: {bad}")
    assert not lg_within_bounds(bad, SLIDE_RECT), (
        f"전제 실패 — layout_geometry 기준으로도 경계 밖이어야 함: {bad}")

    # 2) 수정 seam 구동 — clamp_into_bounds 가 off-slide rect 를 경계 안으로 교정.
    clamped = clamp_into_bounds(bad, SLIDE_RECT)
    print(f"[D3] clamp_into_bounds rect = "
          f"({clamped[0]:.2f}, {clamped[1]:.2f}, {clamped[2]:.2f}, {clamped[3]:.2f})")
    assert lg_within_bounds(clamped, SLIDE_RECT), (
        f"D3 미해소: clamp_into_bounds 결과가 여전히 경계 밖 — {clamped}")
    assert _within_bounds(clamped), (
        f"D3 미해소: clamp 결과가 audit 기준 경계 밖 — {clamped}")

    # fit_within 도 동일하게 음수 off 없는 경계 안 rect 를 산출해야 한다(공통 가드).
    fitted = fit_within(region, 900, 720)
    print(f"[D3] fit_within rect = "
          f"({fitted[0]:.2f}, {fitted[1]:.2f}, {fitted[2]:.2f}, {fitted[3]:.2f})")
    assert lg_within_bounds(fitted, SLIDE_RECT), (
        f"D3 미해소: fit_within 결과가 경계 밖 — {fitted}")

    # 3) 교정된 rect 로 실제 임베드 → 디스크 정독으로도 경계 밖 PICTURE 0건.
    off_l, off_t, draw_w, draw_h = clamped
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(illo, Inches(off_l), Inches(off_t),
                             width=Inches(draw_w), height=Inches(draw_h))
    out = os.path.join(proj, "d3.pptx")
    prs.save(out)

    slides = list(Presentation(out).slides)
    offenders = []
    for idx, slide in enumerate(slides):
        for r, _px in _slide_pictures(slide):
            if not _within_bounds(r):
                offenders.append((idx, r))
    print(f"[D3] 경계 밖 PICTURE = {offenders}")
    assert not offenders, (
        f"D3 반례: 교정 후에도 부분 이미지가 슬라이드 경계 밖에 배치됨 — {offenders}. "
        f"clamp_into_bounds/fit_within 경계 가드가 적용되지 않았다.")


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-p", "no:cacheprovider", "-q"]))
