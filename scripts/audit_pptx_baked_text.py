"""배경 이미지에 '구워진 텍스트(라벨)'가 있는지 정밀 검출 + 이미지 추출.

방법:
  - 풀해상도 그레이스케일에서 '텍스트 줄' 신호 검출:
    각 행(row)의 좌우 밝기 급변(edge) 횟수를 세고, 임계 이상인 행이
    연속/다수면 텍스트 라인이 구워졌을 가능성이 높다.
  - 배경 이미지를 scripts/_audit_out/ 로 추출해 육안 확인 가능하게 저장.
  - 오버레이 텍스트 박스의 위치(인치)도 출력 → 이미지 텍스트와 겹치는 영역 추정.
"""
from __future__ import annotations

import os
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from PIL import Image

EMU_PER_IN = 914400.0
OUT_DIR = os.path.join(os.path.dirname(__file__), "_audit_out")


def _in(emu):
    return round(float(emu) / EMU_PER_IN, 2) if emu is not None else None


def baked_text_score(im: Image.Image):
    """텍스트 줄로 추정되는 행의 비율(%)과 추정 텍스트 줄 수."""
    g = im.convert("L")
    w, h = g.size
    # 가로 1200px 기준으로 정규화(과대 해상도 비용 절감)
    if w > 1200:
        nh = int(h * 1200 / w)
        g = g.resize((1200, max(1, nh)))
        w, h = g.size
    px = list(g.getdata())
    texty_rows = 0
    run = 0
    runs = []
    for y in range(h):
        base = y * w
        trans = 0
        for x in range(w - 1):
            if abs(px[base + x] - px[base + x + 1]) > 45:
                trans += 1
        # 텍스트 줄: 한 행에서 좌우 밝기 급변이 매우 많음(글자 경계)
        if trans >= w * 0.06:
            texty_rows += 1
            run += 1
        else:
            if run >= 3:
                runs.append(run)
            run = 0
    if run >= 3:
        runs.append(run)
    return round(texty_rows / h * 100.0, 1), len(runs)


def main(path):
    os.makedirs(OUT_DIR, exist_ok=True)
    prs = Presentation(path)
    print(f"=== 구워진 텍스트 정밀 검출 + 추출: {path}")
    print(f"추출 경로: {OUT_DIR}\n")
    for idx, slide in enumerate(prs.slides, 1):
        # 오버레이 텍스트 박스 위치
        boxes = []
        for sh in slide.shapes:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                continue
            if getattr(sh, "has_text_frame", False) and (sh.text_frame.text or "").strip():
                boxes.append((sh.name, _in(sh.left), _in(sh.top),
                              _in(sh.width), _in(sh.height),
                              (sh.text_frame.text or "").strip().split("\n")[0][:30]))
        # 배경 이미지
        for sh in slide.shapes:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            try:
                blob = sh.image.blob
                ext = sh.image.ext
            except Exception:
                continue
            fn = os.path.join(OUT_DIR, f"slide{idx:02d}.{ext}")
            with open(fn, "wb") as f:
                f.write(blob)
            try:
                im = Image.open(fn)
                pct, lines = baked_text_score(im)
                baked = pct >= 6.0 or lines >= 6
                print(f"[슬라이드 {idx}] 이미지 {im.size[0]}x{im.size[1]} → "
                      f"텍스트추정행 {pct}%, 텍스트줄 추정 {lines}개  "
                      f"{'⚠ 텍스트 구워짐 강하게 의심' if baked else '구워진 텍스트 적음'}")
            except Exception as e:
                print(f"[슬라이드 {idx}] 이미지 분석 실패: {e}")
        for nm, l, t, w, h, head in boxes:
            print(f"      overlay '{head}'  @({l},{t}) {w}x{h}in")
        print()
    print(f"육안 확인: open {OUT_DIR}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else "")
