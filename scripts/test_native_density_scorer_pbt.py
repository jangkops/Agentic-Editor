"""밀도 채점기 속성 기반 테스트 — spec: pptx-native-density-render (작업 13.1~13.3).

이 파일은 ``scripts/parity_scorer.py`` 의 밀도 채점 경로
(HTML 채점 ``score(html, category)`` + 네이티브 셰이프 채점
``score_native_slide(slide, category)``)에 대한 Correctness Properties
**P9 / P10 / P11** 을 Hypothesis 로 검증한다.

검증 대상 슬라이드는 ``ai_engine/native_layout_renderer.py`` 의 편집가능 네이티브
emit_* (``emit_title``/``emit_section_header_bar``/``emit_contact_box``/
``emit_note_callout``/``emit_numbered_list``/``emit_card_grid``/``emit_figure_slot``)
로 in-memory python-pptx 슬라이드에 직접 방출한다. 베이크_통짜이미지 없이 편집가능
도형만으로 밀도 패리티(reference=6)를 충족함을 확인한다(P9).

각 Property 는 **단일** property-based test 로 구현하고 ``@settings(max_examples>=100)``
이상으로 실행한다. 네트워크 0 — 게이트웨이/Vertex/HTML(Chrome) 렌더를 일절 호출하지
않고, 순수 채점 함수 + in-memory 네이티브 도형 조립만 구동한다.

Run (hermetic):
  ./venv/bin/python -m pytest scripts/test_native_density_scorer_pbt.py -p no:cacheprovider -q

대표 슬라이드 구성 근거(P9):
  - body(8항목, reference 6): section_header(다크 막대) + numbered_item(번호 배지) +
    contact_box(틴트 박스+세로 액센트 바) + note_callout(NOTICE 라벨) +
    figure_slot(회색 플레이스홀더) + slide_footer(하단 텍스트) = 6 검출 → passed.
  - cover(7항목, reference 6): accent_head(대형 제목≥26pt) + step_grid(카드≥2) +
    icon_badge(섹션 헤더 배지 oval) + accent_bar(연락처 세로 바 h≥2.0) +
    notice_chip(노트 콜아웃 박스 h≤0.7) + footer(하단 셰이프) = 6 검출 → passed.
  얇은 단일 레이아웃이 아니라 **충분한 시각요소를 가진 대표 슬라이드**를 구성해
  reference=6 을 채운다(Req 5.1/8.5).

_Requirements: 5.1, 5.2, 5.3, 5.5, 8.5_
_Property: 9, 10, 11_
"""
from __future__ import annotations

import os
import sys

# ai_engine(repo 루트)와 채점 도구(scripts/)를 import 가능하게 한다.
_HERE = os.path.dirname(os.path.abspath(__file__))
_ROOT = os.path.dirname(_HERE)
for _p in (_ROOT, _HERE):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from hypothesis import given, settings, strategies as st, HealthCheck  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.util import Inches  # noqa: E402

# 채점기(테스트 대상) — 순수 함수, 네트워크 0.
from parity_scorer import (  # noqa: E402
    score,
    score_native_slide,
    COVER_CHECKLIST,
    BODY_CHECKLIST,
)

# 네이티브 emit_* (편집가능 도형 방출) + 디자인 토큰 — 네트워크 0.
from ai_engine.native_layout_renderer import (  # noqa: E402
    emit_title,
    emit_section_header_bar,
    emit_contact_box,
    emit_note_callout,
    emit_numbered_list,
    emit_card_grid,
    emit_figure_slot,
    _emit_text_block,
)
from ai_engine.slide_templates import design_tokens_for_profile  # noqa: E402

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5


# ===========================================================================
# 공통 헬퍼 — in-memory 16:9 슬라이드
# ===========================================================================
def _blank_169_presentation():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    return prs


def _blank_slide(prs):
    # layout[6] = 완전 빈 레이아웃(placeholder 간섭 없음).
    return prs.slides.add_slide(prs.slide_layouts[6])


def _tokens():
    # 기본 프로파일 토큰(SLIDE_DESIGN 사본) — 신규 토큰 정의 없음.
    return design_tokens_for_profile(None)


# ===========================================================================
# 텍스트 생성기 — 한글 / 특수문자 / ascii, 항상 비어있지 않음(strip 후 폴백)
# ===========================================================================
_KO = st.sampled_from([
    "프로젝트 개요", "핵심 전략", "분기 성과", "로드맵", "리스크 분석",
    "팀 구성", "목표 정의", "요약 노트", "연락처", "주요 지표",
])
_SPECIAL = st.sampled_from([
    "A/B 테스트", "R&D 100%", "#1 우선순위", "α→β 전환", "—강조—",
    "C++ / Rust", "KPI: 95%", "<태그>&값", '"인용"',
])
_ASCII = st.text(
    alphabet=st.characters(min_codepoint=32, max_codepoint=126),
    min_size=1, max_size=30,
)


def _clean(s: str) -> str:
    """공백 정리 후 비면 한국어 폴백 — 콘텐츠 텍스트는 항상 비어있지 않게."""
    return (s or "").strip() or "내용"


text_strategy = st.one_of(_KO, _SPECIAL, _ASCII).map(_clean)
lines_strategy = st.lists(text_strategy, min_size=1, max_size=5)

# 대표 body 슬라이드 입력(충분한 시각요소).
body_spec = st.fixed_dictionaries({
    "header_no": st.integers(min_value=1, max_value=9),
    "header_title": text_strategy,
    "contact": st.lists(text_strategy, min_size=1, max_size=3),
    "note": text_strategy,
    "numbered": st.lists(text_strategy, min_size=1, max_size=5),
    "footer": text_strategy,
})

# 대표 cover 슬라이드 입력(충분한 시각요소).
cover_spec = st.fixed_dictionaries({
    "title": text_strategy,
    "cards": st.lists(text_strategy, min_size=2, max_size=4),
    "header_no": st.integers(min_value=1, max_value=9),
    "header_title": text_strategy,
    "contact": st.lists(text_strategy, min_size=1, max_size=3),
    "note": text_strategy,
    "footer": text_strategy,
})


# ===========================================================================
# 네이티브 슬라이드 빌더 — emit_* 로 편집가능 도형 방출
# enabled 집합으로 요소 일부만 방출(부분 슬라이드 생성: P10 가변 입력용).
# ===========================================================================
_BODY_ELEMENTS = ("header", "contact", "note", "numbered", "figure", "footer")
_COVER_ELEMENTS = ("title", "cards", "header", "contact", "note", "footer")


def _build_body_slide(slide, tokens, spec, enabled=None):
    """body 대표 슬라이드 — section_header / numbered / contact / note / figure /
    footer 를 편집가능 네이티브 도형으로 방출한다(검출기 기준 충족)."""
    enabled = set(enabled if enabled is not None else _BODY_ELEMENTS)

    if "header" in enabled:
        # 다크 막대 + 번호 배지(digit) + 제목 → section_header & numbered_item.
        emit_section_header_bar(
            slide, spec["header_no"], spec["header_title"], tokens,
            (0.7, 0.4, 11.9, 1.0),
        )
    if "contact" in enabled:
        # 틴트 박스(border) + 세로 액센트 바(w0.12,h1.5) → contact_box.
        emit_contact_box(slide, spec["contact"], tokens, (0.7, 1.6, 5.5, 1.5))
    if "note" in enabled:
        # NOTICE 라벨 → note_callout.
        emit_note_callout(slide, spec["note"], tokens, (6.6, 1.6, 6.0, 1.5))
    if "numbered" in enabled:
        # 번호 배지(oval, digit) → numbered_item.
        emit_numbered_list(slide, spec["numbered"], tokens, (0.7, 3.3, 5.5, 2.5))
    if "figure" in enabled:
        # 회색 플레이스홀더(텍스트 없음) → figure_slot.
        emit_figure_slot(slide, (6.6, 3.3, 5.5, 2.5))
    if "footer" in enabled:
        # 하단 텍스트(top≥6.5, h≤0.7) → slide_footer.
        _emit_text_block(slide, spec["footer"], tokens,
                         (0.7, 6.8, 6.0, 0.4), role="caption")


def _build_cover_slide(slide, tokens, spec, enabled=None):
    """cover 대표 슬라이드 — accent_head / step_grid / icon_badge / accent_bar /
    notice_chip / footer 를 편집가능 네이티브 도형으로 방출한다(검출기 기준 충족)."""
    enabled = set(enabled if enabled is not None else _COVER_ELEMENTS)

    if "title" in enabled:
        # 대형 제목(title role 30pt ≥26) → accent_head.
        emit_title(slide, spec["title"], tokens, (0.7, 0.4, 11.9, 1.4))
    if "cards" in enabled:
        # 라운드 카드 ≥2 (w≥1.2, h≥0.8) → step_grid.
        emit_card_grid(slide, list(spec["cards"]), tokens, (0.7, 2.0, 9.0, 2.0))
    if "header" in enabled:
        # 섹션 헤더 배지(oval accent, 소형) → icon_badge.
        emit_section_header_bar(
            slide, spec["header_no"], spec["header_title"], tokens,
            (0.7, 4.3, 9.0, 1.0),
        )
    if "contact" in enabled:
        # 세로 액센트 바(w0.12, h4.5≥2.0) → accent_bar.
        emit_contact_box(slide, spec["contact"], tokens, (10.0, 0.4, 3.0, 4.5))
    if "note" in enabled:
        # 노트 콜아웃 박스(라운드, h0.6≤0.7) → notice_chip.
        emit_note_callout(slide, spec["note"], tokens, (0.7, 6.5, 5.0, 0.6))
    if "footer" in enabled:
        # 하단 텍스트(top≥6.3) → footer.
        _emit_text_block(slide, spec["footer"], tokens,
                         (0.7, 6.7, 4.0, 0.4), role="caption")


def _build_native_slide(category, spec, enabled=None):
    prs = _blank_169_presentation()
    slide = _blank_slide(prs)
    tokens = _tokens()
    if category == "cover":
        _build_cover_slide(slide, tokens, spec, enabled=enabled)
    else:
        _build_body_slide(slide, tokens, spec, enabled=enabled)
    return slide


# ===========================================================================
# 채점 결과 내부 정합 단언(공통) — items/missing/density/passed 의 무모순성.
# ===========================================================================
def _assert_report_consistent(res):
    items = res["items"]
    present_names = [it["name"] for it in items if it["present"]]
    missing_names = [it["name"] for it in items if not it["present"]]

    # total 은 items 개수와 일치.
    assert res["total"] == len(items)
    # density_score 는 present 항목 수와 정확히 일치.
    assert res["density_score"] == len(present_names)
    # missing 목록은 실제 미검출 항목과 정확히 일치(순서 포함).
    assert res["missing"] == missing_names
    # present ∪ missing == 전체, 상호배타.
    assert set(present_names).isdisjoint(set(missing_names))
    assert set(present_names) | set(missing_names) == {it["name"] for it in items}
    # 합격 판정은 density >= reference 와 동치.
    assert res["passed"] == (res["density_score"] >= res["reference_score"])


# ===========================================================================
# Property 9 — 네이티브만으로 밀도 패리티 합격
# ===========================================================================
# Feature: pptx-native-density-render, Property 9: 알려진_레이아웃으로 네이티브 렌더된
# 슬라이드(카테고리 cover 또는 body)에 대해, score_native_slide 의 density_score 는 해당
# 카테고리의 reference_score(cover=6, body=6) 이상이며 passed=True 이다 — 베이크_통짜이미지
# 없이 편집가능 셰이프만으로 합격한다. (Validates: Requirements 5.1, 8.5)
@settings(max_examples=120, deadline=None,
          suppress_health_check=[HealthCheck.too_slow])
@given(category=st.sampled_from(["cover", "body"]), data=st.data())
def test_property9_native_density_parity_pass(category, data):
    spec = data.draw(cover_spec if category == "cover" else body_spec)
    slide = _build_native_slide(category, spec)

    res = score_native_slide(slide, category)

    # 편집가능 네이티브 도형만으로 reference(6) 이상 + 합격.
    assert res["reference_score"] == 6
    assert res["density_score"] >= res["reference_score"], (
        f"{category}: density {res['density_score']} < reference "
        f"{res['reference_score']}, missing={res['missing']}"
    )
    assert res["passed"] is True
    assert res["missing"] == [] or res["density_score"] >= res["reference_score"]
    # 내부 정합도 함께 보장.
    _assert_report_consistent(res)


# ===========================================================================
# Property 10 — 밀도 채점기의 결정성과 보고 완전성
# ===========================================================================
def _html_from_subset(category, indices):
    """선택된 체크리스트 인덱스의 마커만 포함하는 HTML 문자열을 만든다.
    score(html, category) 는 marker 부분문자열 존재로 검출하므로, 포함한 마커만
    present 가 된다(결정적)."""
    checklist = COVER_CHECKLIST if category == "cover" else BODY_CHECKLIST
    chosen = sorted(set(indices))
    spans = [f'<div {checklist[i][1]}></div>' for i in chosen]
    return "<section>" + "".join(spans) + "</section>"


# Feature: pptx-native-density-render, Property 10: 유효한 (입력, 카테고리) 쌍에 대해,
# score/score_native_slide 는 passed == (density_score >= reference_score) 를 보장하고,
# passed=False 일 때 missing 목록이 비어있지 않으며 실제 누락된 체크리스트 항목과 일치한다.
# (Validates: Requirements 5.2, 5.3)
@settings(max_examples=150, deadline=None,
          suppress_health_check=[HealthCheck.too_slow])
@given(category=st.sampled_from(["cover", "body"]),
       kind=st.sampled_from(["native", "html"]),
       data=st.data())
def test_property10_determinism_and_report_completeness(category, kind, data):
    if kind == "html":
        total = len(COVER_CHECKLIST if category == "cover" else BODY_CHECKLIST)
        # 0..total 개의 임의 마커 부분집합 → 합격/불합격 모두 분포.
        indices = data.draw(st.lists(st.integers(0, total - 1),
                                     min_size=0, max_size=total))
        html = _html_from_subset(category, indices)
        res = score(html, category)
    else:
        # 임의 요소 부분집합으로 네이티브 슬라이드 구성(합격/불합격 모두 분포).
        elements = _COVER_ELEMENTS if category == "cover" else _BODY_ELEMENTS
        enabled = data.draw(st.lists(st.sampled_from(elements),
                                     min_size=0, max_size=len(elements)).map(set))
        spec = data.draw(cover_spec if category == "cover" else body_spec)
        slide = _build_native_slide(category, spec, enabled=enabled)
        res = score_native_slide(slide, category)

    # 결정성/보고 완전성 — 임의 입력에서 항상 성립.
    _assert_report_consistent(res)

    # passed == (density >= reference) (명시 단언).
    assert res["passed"] == (res["density_score"] >= res["reference_score"])

    # 불합격이면 missing 비어있지 않고 실제 누락 항목과 일치.
    if not res["passed"]:
        assert len(res["missing"]) > 0
        expected_missing = [it["name"] for it in res["items"] if not it["present"]]
        assert res["missing"] == expected_missing


# ===========================================================================
# Property 11 — 빈/미지원 입력은 ValueError 를 발생시킨다
# ===========================================================================
_VALID_CATEGORIES = ("cover", "body")
# 미지원 카테고리(cover·body 외) — 대소문자/공백/유사어 포함.
_BAD_CATEGORIES = st.sampled_from([
    "timeline", "section_divider", "table", "Cover", "BODY", "body ",
    " cover", "", "two_column", "feature_grid", None, "slide",
])


# Feature: pptx-native-density-render, Property 11: 빈 입력(None 또는 빈 문자열) 또는
# 미지원 카테고리(cover·body 외)에 대해, score/score_native_slide 는 점수를 산출하지 않고
# ValueError 를 발생시킨다. (Validates: Requirements 5.5)
@settings(max_examples=150, deadline=None,
          suppress_health_check=[HealthCheck.too_slow])
@given(func=st.sampled_from(["score", "native"]),
       scenario=st.sampled_from(["bad_category", "empty_input"]),
       data=st.data())
def test_property11_empty_or_unsupported_raises_value_error(func, scenario, data):
    if func == "score":
        if scenario == "bad_category":
            category = data.draw(_BAD_CATEGORIES)
            # 미지원 카테고리는 입력과 무관하게 ValueError(카테고리 먼저 검사).
            html = data.draw(st.one_of(st.none(), st.just(""),
                                       text_strategy.map(lambda s: f"<div>{s}</div>")))
            with pytest.raises(ValueError):
                score(html, category)
        else:  # empty_input — 유효 카테고리 + 빈 입력(None/"").
            category = data.draw(st.sampled_from(_VALID_CATEGORIES))
            html = data.draw(st.sampled_from([None, ""]))
            with pytest.raises(ValueError):
                score(html, category)
    else:  # score_native_slide
        if scenario == "bad_category":
            category = data.draw(_BAD_CATEGORIES)
            # 카테고리가 먼저 검사되므로 slide 가 None 이어도 ValueError.
            slide = data.draw(st.sampled_from(["__build__", None]))
            real_slide = None
            if slide == "__build__":
                real_slide = _build_native_slide(
                    "body", data.draw(body_spec))
            with pytest.raises(ValueError):
                score_native_slide(real_slide, category)
        else:  # empty_input — 유효 카테고리 + slide=None.
            category = data.draw(st.sampled_from(_VALID_CATEGORIES))
            with pytest.raises(ValueError):
                score_native_slide(None, category)
