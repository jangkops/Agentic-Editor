#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Regression — build_native_diagram의 NOTICE 콜아웃(note 파라미터).

doc2(온보딩 매뉴얼) 시그니처 요소: 크림 배경 + 주황 좌측 보더 + 'NOTICE' 라벨 + 본문.
모두 편집 가능한 네이티브 도형이어야 한다(통짜 이미지 금지).

실행: python scripts/test_native_notice_callout.py
"""
import sys
from pathlib import Path

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import ai_engine.native_diagram_pptx as nd


def _fill_hex(sh):
    try:
        return str(sh.fill.fore_color.rgb)
    except Exception:
        return "?"


def _new_slide():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def test_callout_renders():
    prs, s = _new_slide()
    ok = nd.build_native_diagram(
        s, "cards",
        "성능 최적화: 빠른 응답\n클라우드 인프라: AWS 확장성\n보안 강화: 권한 통제",
        title="핵심 항목",
        note="미설치 시 사내망 접속이 차단됩니다. 반드시 3종 모두 설치하세요.",
    )
    assert ok, "build_native_diagram returned False"
    autos = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    pics = [sh for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE]
    cream = [sh for sh in autos if _fill_hex(sh) == "FFF8EE"]
    orange = [sh for sh in autos if _fill_hex(sh) == "ED7D31"]
    notice = any(sh.has_text_frame and "NOTICE" in sh.text_frame.text for sh in s.shapes)
    body = any(sh.has_text_frame and "사내망 접속이 차단" in sh.text_frame.text for sh in s.shapes)
    # 풀블리드 통짜 이미지 없음(편집 가능 불변식)
    fb = [p for p in pics if (p.width or 0) >= int(13 * 914400)]
    assert len(cream) == 1, f"cream box={len(cream)}"
    assert len(orange) >= 1, f"orange border={len(orange)}"
    assert notice, "NOTICE label missing"
    assert body, "note body missing"
    assert not fb, "full-bleed raster present"
    print(f"[callout] ok cream={len(cream)} orange={len(orange)} NOTICE={notice} body={body}")


def test_no_note_no_callout():
    prs, s = _new_slide()
    ok = nd.build_native_diagram(
        s, "cards",
        "성능 최적화: 빠른 응답\n보안 강화: 권한 통제",
        title="핵심 항목",
    )
    assert ok
    cream = [sh for sh in s.shapes
             if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and _fill_hex(sh) == "FFF8EE"]
    assert len(cream) == 0, "callout drawn without note"
    notice = any(sh.has_text_frame and "NOTICE" in sh.text_frame.text for sh in s.shapes)
    assert not notice, "NOTICE label drawn without note"
    print("[no-note] ok — 콜아웃 없음")


if __name__ == "__main__":
    test_callout_renders()
    test_no_note_no_callout()
    print("ALL PASS")
