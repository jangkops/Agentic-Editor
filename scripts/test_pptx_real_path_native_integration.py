"""실제 생성 경로 네이티브 통합 게이트 — pptx-native-density-render (작업 20).

이 테스트는 **실제 ``ai_engine.server._tool_generate_pptx`` 를 구동**하여, 옵트인 없이
(기본 동작으로) 표지·본문·다이어그램 슬라이드가 전부 **편집가능 네이티브 + 겹침<10% +
경계 안 + 제목 1회** 로 생성됨을 산출물 audit 으로 단언한다.

기존 ``test_pptx_native_density_audit_integration.py`` 는 ``render_native_layout`` 을
*직접* 호출해 덱을 만들었다(라이브러리 단위 검증). 본 테스트는 그와 달리 **실제 PPTX
생성 함수(_tool_generate_pptx)의 전체 제어흐름**을 통과시켜, 근본 원인(옵트인 OFF,
레이아웃 불일치, 표지/다이어그램 겹침)이 실제 경로에서 해소됐음을 입증한다.

헤르메틱(네트워크 0 · Chrome 미사용):
  - 게이트웨이 LLM(``_llm_pick_slide_layout``/``_get_gw``/``_specialized_model_for_task``),
    Vertex(옵트인 OFF + ``AE_PREFER_VERTEX_IMAGE=0``), Chrome/브리지
    (``_render_html_slide_to_png``/``_find_local_chrome``/``_call_bridge``)를 전부 mock.
  - ``_render_html_slide_to_png`` mock 은 **콘텐츠가 구워진(baked) PNG** 를 반환하도록
    구성 — 통합이 그 베이크 배경을 콘텐츠로 채택하지 않고 네이티브로 렌더함을 입증한다
    (회귀 방지: 통짜 이미지로 되돌아가면 (a)편집텍스트/(f)베이크 검사에서 걸린다).
  - 20초+ hang 방지: 실제 subprocess/네트워크 없음(전부 in-memory mock).

실행:
  ./venv/bin/python -m pytest scripts/test_pptx_real_path_native_integration.py \
      -p no:cacheprovider -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
import tempfile
import shutil
import contextlib
from unittest.mock import patch

_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
_REPO_ROOT = os.path.dirname(_THIS_DIR)
_AE_DIR = os.path.join(_REPO_ROOT, "ai_engine")
for _p in (_REPO_ROOT, _AE_DIR, _THIS_DIR):
    if _p not in sys.path:
        sys.path.insert(0, _p)

import pytest  # noqa: E402
from PIL import Image, ImageDraw  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: E402

import server  # noqa: E402  (ai_engine on path)
from slide_templates import design_tokens_for_profile  # noqa: E402
from audit_pptx_native_density import (  # noqa: E402  (scripts on path)
    audit_native_density,
    AuditReport,
    CHECK_EDITABLE_TEXT,
    CHECK_OVERLAP,
    CHECK_BOUNDS,
    CHECK_TITLE_ONCE,
    CHECK_BAKED_TEXT,
    CHECK_ZORDER,
)

_TOKENS = design_tokens_for_profile(None)

# 실제 경로가 네이티브로 감을 입증하는 핵심 검사 집합:
#   (a) 편집가능 텍스트 런 ≥1, (b) 겹침<10%, (c) 경계 안, (d) 제목 1회,
#   (f) 풀블리드 베이크 텍스트 미검출, (g) 텍스트 z-order 이미지 위.
# 밀도(e)/스타일(h)는 렌더러별 시각 완성도 차원으로 별도 보고(핵심 게이트와 직교).
_CORE_CHECKS = [
    CHECK_EDITABLE_TEXT, CHECK_OVERLAP, CHECK_BOUNDS,
    CHECK_TITLE_ONCE, CHECK_BAKED_TEXT, CHECK_ZORDER,
]


# ===========================================================================
# 본문 슬라이드별 LLM 레이아웃 픽(mock) — 실제 게이트웨이 호출 0
# ===========================================================================
# 헤딩별로 (layout, data) 를 반환한다. 특히 "핵심 지표" 슬라이드는 **비-7키**
# ``kpi_summary`` 를 반환해, 서버의 map_to_native_layout(작업 B)이 feature_grid 로
# 매핑해 네이티브 라우팅됨을 입증한다.
_PICK_TABLE = {
    "구축 로드맵": ("timeline", {
        "title": "구축 로드맵",
        "subtitle": "4단계 추진 일정",
        "steps": [
            {"label": "1분기", "title": "기반 구축", "description": "네이티브 렌더러 스캐폴드"},
            {"label": "2분기", "title": "요소 매핑", "description": "고밀도 요소 도형화"},
            {"label": "3분기", "title": "검증 게이트", "description": "산출물 audit 구축"},
            {"label": "4분기", "title": "정식 출시", "description": "전사 배포 및 안정화"},
        ],
    }),
    "시스템 아키텍처": ("architecture", {
        "title": "시스템 아키텍처",
        "subtitle": "레이어 구성",
        "layers": [
            {"name": "프레젠테이션", "description": "네이티브 도형 렌더", "items": ["emit_*", "tokens"]},
            {"name": "기하 보정", "description": "겹침/경계 보정", "items": ["resolve_collisions"]},
            {"name": "검증", "description": "산출물 audit", "items": ["audit_native_density"]},
        ],
    }),
    "현황과 과제": ("two_column", {
        "title": "현황과 과제",
        "subtitle": "내부 역량과 외부 환경",
        "left_content": ["강력한 데이터 자산", "검증된 추론 파이프라인", "안정적 인프라"],
        "right_content": ["멀티모달 대응 지연", "추론 비용 최적화", "거버넌스 보강"],
    }),
    "핵심 지표": ("kpi_summary", {  # 비-7키 → map_to_native_layout → feature_grid
        "title": "핵심 지표",
        "subtitle": "정량 성과",
        "features": [
            {"icon": "★", "title": "밀도", "description": "젠스파크급 고밀도 유지"},
            {"icon": "◆", "title": "편집성", "description": "모든 텍스트 편집 가능"},
            {"icon": "●", "title": "겹침", "description": "10% 미만 보장"},
            {"icon": "▲", "title": "장식", "description": "z-order 분리"},
        ],
    }),
}


async def _fake_pick(gw, model_id, section_heading, section_body,
                     doc_context="", role="", bullet_count=0):
    key = (section_heading or "").strip()
    if key in _PICK_TABLE:
        layout, data = _PICK_TABLE[key]
        return {"layout": layout, "data": dict(data)}
    # 미지정 헤딩 → two_column 기본(콘텐츠 텍스트 보존)
    return {"layout": "two_column", "data": {
        "title": key or "내용", "subtitle": "",
        "left_content": [l for l in (section_body or "").splitlines() if l.strip()][:3] or ["항목"],
        "right_content": ["세부 사항", "추가 근거"],
    }}


def _baked_png(path: str):
    """콘텐츠가 '구워진' 것처럼 보이는 PNG(수평 텍스트 줄 다수) — 채택되면 (f) 실패.

    통합이 이 베이크 배경을 콘텐츠로 채택하지 않고 네이티브로 렌더함을 입증하기 위한
    함정 이미지다.
    """
    img = Image.new("RGB", (1920, 1080), (245, 246, 248))
    d = ImageDraw.Draw(img)
    # 텍스트 줄을 흉내내는 어두운 수평 막대 다수(baked_text 판정 유도).
    y = 120
    for _ in range(14):
        d.rectangle([160, y, 1500, y + 26], fill=(30, 33, 40))
        y += 60
    img.save(path, "PNG")


async def _fake_render_html(html, output_path, width=1920, height=1080, timeout=30):
    """Chrome 브리지 mock — 콘텐츠-베이크 PNG 를 디스크에 쓰고 성공 envelope 반환."""
    try:
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        _baked_png(output_path)
        return {"ok": True, "path": output_path, "width": width, "height": height,
                "sizeBytes": os.path.getsize(output_path)}
    except Exception as e:  # pragma: no cover
        return {"ok": False, "error": str(e)}


async def _fake_html_section(*a, **k):
    """본문 HTML 섹션 렌더 mock — 빈 문자열(네이티브 라우팅 시 미사용). 회귀 안전."""
    return ""


@contextlib.contextmanager
def _hermetic_env():
    """헤르메틱 환경 — Vertex/HTML 강제 상태 고정. 원복 보장."""
    keys = {
        "AE_ENABLE_HTML_SLIDES": "1",     # HTML 경로 ON(→ 네이티브 라우팅 게이트 전제)
        "AE_DISABLE_HTML_SLIDES": "0",
        "AE_PREFER_VERTEX_IMAGE": "0",    # Vertex 이미지 병렬 생성 차단(네트워크 0)
        "AE_ENABLE_VERTEX_IMAGE": "0",    # 네이티브 장식 옵트인 OFF
        "AE_ENABLE_VERTEX_BG": "0",       # 공유 배경 OFF
        # AE_NATIVE_LAYOUT_RENDER 는 설정하지 않음 → 기본 활성(작업 A) 검증.
    }
    old = {k: os.environ.get(k) for k in keys}
    os.environ.pop("AE_NATIVE_LAYOUT_RENDER", None)
    os.environ.update(keys)
    try:
        yield
    finally:
        for k, v in old.items():
            if v is None:
                os.environ.pop(k, None)
            else:
                os.environ[k] = v


@contextlib.contextmanager
def _all_mocks():
    """게이트웨이/Chrome/모델 선택을 전부 mock (네트워크 0)."""
    with patch.object(server, "_llm_pick_slide_layout", _fake_pick), \
         patch.object(server, "_render_html_slide_to_png", _fake_render_html), \
         patch.object(server, "_generate_html_slide_for_section", _fake_html_section), \
         patch.object(server, "_find_local_chrome", lambda *a, **k: "/fake/chrome"), \
         patch.object(server, "_call_bridge", lambda *a, **k: None), \
         patch.object(server, "_get_gw", lambda *a, **k: object()), \
         patch.object(server, "_specialized_model_for_task", lambda *a, **k: "mock-model"):
        yield


def _tool_input():
    """표지 + 본문 4장(다이어그램 2 + 2단 + kpi) 요청."""
    return {
        "title": "차세대 AI 플랫폼 통합 전략 보고서",
        "subtitle": "편집가능 네이티브 렌더링",
        "eyebrow": "2026 전략",
        "slides": [
            {"title": "구축 로드맵", "bullets": ["1분기 기반", "2분기 매핑", "3분기 검증", "4분기 출시"]},
            {"title": "시스템 아키텍처", "bullets": ["프레젠테이션", "기하 보정", "검증 레이어"]},
            {"title": "현황과 과제", "bullets": ["데이터 자산", "추론 파이프라인", "인프라"]},
            {"title": "핵심 지표", "bullets": ["밀도", "편집성", "겹침", "장식"]},
        ],
    }


def _generate_pptx_real_path(project_path: str) -> str:
    """실제 _tool_generate_pptx 를 구동하고 저장된 .pptx 절대경로를 반환."""
    with _hermetic_env(), _all_mocks():
        result_str = asyncio.run(server._tool_generate_pptx(
            _tool_input(), project_path, aws_profile="", bedrock_user=""))
    result = json.loads(result_str)
    assert "error" not in result, f"생성 실패: {result}"
    abspath = result.get("absPath") or ""
    assert abspath and os.path.isfile(abspath), f"산출물 없음: {result}"
    return abspath


def _count_editable_text_shapes(slide) -> int:
    n = 0
    for sh in slide.shapes:
        try:
            if sh.has_text_frame:
                txt = "".join(r.text for p in sh.text_frame.paragraphs for r in p.runs)
                if txt.strip():
                    n += 1
        except Exception:
            pass
    return n


def _count_pictures(slide) -> int:
    n = 0
    for sh in slide.shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                n += 1
        except Exception:
            pass
    return n


def _fail_summary(report: AuditReport) -> str:
    return "; ".join(
        f"[slide {f['slide']}] {f['check']}: {f['signal']}" for f in report.failures
    ) or "(없음)"


# ===========================================================================
# 게이트 1 — 실제 경로: 모든 슬라이드 편집가능 + 겹침<10% + 경계 안 + 제목 1회
# ===========================================================================

def test_real_path_all_slides_native_editable_no_overlap():
    """실제 _tool_generate_pptx 구동(옵트인 없음) → 표지+본문 전부:
      (a) 편집가능 텍스트 ≥1, (b) 겹침<10%, (c) 경계 안, (d) 제목 1회,
      (f) 베이크 텍스트 미검출, (g) z-order 정상 — 핵심 검사 전수 통과.

    **Validates: Requirements 1.1, 1.2, 1.5, 2.1, 3.1, 4.1, 7.1, 7.2, 7.3, 8.1, 8.2, 8.3, 8.4, 8.8, 8.9**
    """
    project = tempfile.mkdtemp(prefix="rp_native_")
    try:
        pptx_path = _generate_pptx_real_path(project)

        # 슬라이드 수 = 표지 + 본문 4 = 5
        prs = Presentation(pptx_path)
        slides = list(prs.slides)
        assert len(slides) == 5, f"슬라이드 수 예상 5, 실제 {len(slides)}"

        # 모든 슬라이드에 편집가능 텍스트가 존재(통짜 이미지 아님).
        for idx, slide in enumerate(slides):
            assert _count_editable_text_shapes(slide) >= 1, (
                f"슬라이드 {idx+1}: 편집가능 텍스트 0 (통짜 이미지 의심)")

        # 산출물 audit — 핵심 검사(a,b,c,d,f,g) 전수 실패 없음.
        report = audit_native_density(pptx_path, _TOKENS)
        assert isinstance(report, AuditReport)
        failed = {f["check"] for f in report.failures}
        for chk in _CORE_CHECKS:
            assert chk not in failed, (
                f"핵심 검사 실패: {chk} — {_fail_summary(report)}")
    finally:
        shutil.rmtree(project, ignore_errors=True)


# ===========================================================================
# 게이트 2 — 회귀 방지: 콘텐츠-베이크 배경이 채택되지 않음(네이티브 우선)
# ===========================================================================

def test_real_path_does_not_adopt_baked_fullbleed():
    """_render_html_slide_to_png mock 이 콘텐츠-베이크 PNG 를 반환해도, 통합이 그
    베이크 배경을 콘텐츠로 채택하지 않음을 (f)베이크 검사 통과 + 편집텍스트 존재로
    입증한다(회귀 방지: 통짜 이미지 회귀 시 여기서 실패).

    **Validates: Requirements 1.5, 6.1, 6.3, 8.8**
    """
    project = tempfile.mkdtemp(prefix="rp_baked_")
    try:
        pptx_path = _generate_pptx_real_path(project)
        report = audit_native_density(pptx_path, _TOKENS)
        failed = {f["check"] for f in report.failures}
        # 베이크 텍스트가 풀블리드로 채택됐다면 (f) 가 실패한다.
        assert CHECK_BAKED_TEXT not in failed, (
            f"콘텐츠-베이크 배경이 채택됨(통짜 이미지 회귀): {_fail_summary(report)}")

        # 표지 포함 모든 슬라이드에 편집가능 콘텐츠 텍스트 공존.
        prs = Presentation(pptx_path)
        for idx, slide in enumerate(prs.slides):
            assert _count_editable_text_shapes(slide) >= 1, (
                f"슬라이드 {idx+1}: 편집가능 텍스트 없음")
    finally:
        shutil.rmtree(project, ignore_errors=True)


# ===========================================================================
# 게이트 3 — 표지: 네이티브 편집 표지 + 겹침/제목 규칙(콘텐츠-베이크 미채택)
# ===========================================================================

def test_real_path_cover_is_native_editable():
    """표지(첫 슬라이드)가 편집가능 네이티브 텍스트를 갖고, (b)겹침/(d)제목 1회
    규칙을 만족함을 단언(작업 C — 네이티브 표지 우선, 콘텐츠-베이크 HTML 표지 미채택).

    **Validates: Requirements 4.1, 4.3, 7.1**
    """
    project = tempfile.mkdtemp(prefix="rp_cover_")
    try:
        pptx_path = _generate_pptx_real_path(project)
        prs = Presentation(pptx_path)
        cover = list(prs.slides)[0]
        assert _count_editable_text_shapes(cover) >= 1, "표지에 편집가능 텍스트 없음"

        report = audit_native_density(pptx_path, _TOKENS)
        cover_fail = {f["check"] for f in report.failures if f["slide"] == 1}
        for chk in (CHECK_OVERLAP, CHECK_BOUNDS, CHECK_TITLE_ONCE,
                    CHECK_EDITABLE_TEXT, CHECK_BAKED_TEXT):
            assert chk not in cover_fail, (
                f"표지 검사 실패: {chk} — {_fail_summary(report)}")
    finally:
        shutil.rmtree(project, ignore_errors=True)
