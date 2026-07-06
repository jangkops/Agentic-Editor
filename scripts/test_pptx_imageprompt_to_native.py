"""Regression — 구조/다이어그램 슬라이드가 imagePrompt만 있어도 편집 가능한 네이티브 도형이 된다.

사용자 보고(스크린샷): 슬라이드 내부에 폴더 트리/아이소메트릭 박스 *래스터 이미지*가
생성됨 → 내부 요소(폴더/화살표/라벨)를 하나하나 수정할 수 없음.

근본 원인: 모델이 구조 슬라이드를 본문(bullets) 없이 imagePrompt 문장만으로 만들면,
`_classify_section_diagram`이 (a) 영어 키워드(structure/repository/hierarchy 등)를 못 잡거나
(b) 잡아도 구조화된 노드를 못 만들어 네이티브 빌더가 실패 → 통짜 래스터로 폴백.

수정:
  - 키워드 확장(structure/repository/hierarchy/구성도 등) + heading뿐 아니라 본문/프롬프트도 검사.
  - `_mine_diagram_entities`로 프롬프트 prose에서 노드 라벨(폴더명 등)을 추출 →
    루트+자식 트리/블록으로 네이티브 도형 생성.

Correctness properties:
  P1. 콤마/슬래시/화살표 prose에서 엔티티가 정확히 추출된다(설명 불용어 제외).
  P2. "folder tree ... frontend, backend, tests, docs, scripts" 프롬프트 → tree로 분류 + 다수 노드.
  P3. imagePrompt만 있는 구조 슬라이드로 만든 PPTX는 autoshape>0 + Picture==0(편집 가능).
  P4. 다이어그램 신호가 없는 순수 사진 프롬프트는 네이티브로 바뀌지 않는다(빈 분류).

실행: pytest scripts/test_pptx_imageprompt_to_native.py -q
"""
from __future__ import annotations

import os
import sys
import json
import asyncio
from pathlib import Path

import pytest

_ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(_ROOT))
sys.path.insert(0, str(_ROOT / "ai_engine"))

server = pytest.importorskip("ai_engine.server")
pptx = pytest.importorskip("pptx")


def test_mine_entities_basic():
    """P1 — 구분자/불용어 처리."""
    ents = server._mine_diagram_entities(
        "folder tree of repository with frontend, backend, tests, docs, scripts")
    assert "frontend" in ents and "backend" in ents and "scripts" in ents
    # 설명 불용어는 제거
    assert "folder" not in ents and "repository" not in ents and "tree" not in ents


def test_mine_entities_slash_and_arrow():
    ents = server._mine_diagram_entities("user -> gateway -> backend / database")
    assert ents == ["user", "gateway", "backend", "database"]


def test_classify_imageprompt_structure_to_tree():
    """P2 — 본문 없이 구조 프롬프트만 → tree + 다수 노드."""
    kind, content = server._classify_section_diagram(
        "프로젝트 구조", "", "deck")
    # heading '구조' → tree, 본문 없음 → 빈 분류일 수 있으나 프롬프트 경로로 보강됨
    # 본문 없이 heading만이면 마이닝 소스가 없어 빈 분류 → 프롬프트 재시도 경로를 직접 검증
    kind2, content2 = server._classify_section_diagram(
        "프로젝트 구조",
        "folder tree of repository with frontend, backend, tests, docs, scripts")
    assert kind2 == "tree", f"tree로 분류되지 않음: {kind2}"
    nodes = [ln for ln in content2.splitlines() if ln.strip()]
    assert len(nodes) >= 5, f"노드 부족: {nodes}"


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


def _analyze(pptx_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(pptx_path)
    out = {"autoshapes": 0, "pictures": 0}
    for slide in prs.slides:
        for shp in slide.shapes:
            st = shp.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                out["pictures"] += 1
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                out["autoshapes"] += 1
    return out


def test_imageprompt_structure_slide_becomes_native(tmp_path, monkeypatch):
    """P3 — imagePrompt만 있는 구조 슬라이드 → 네이티브 도형, 래스터 미생성."""
    os.environ["AE_GENERATED_ROOT"] = str(tmp_path)

    called = {"n": 0}

    async def _fake_img(tool_input, project_path, aws_profile='', bedrock_user=''):
        called["n"] += 1  # 호출되면(래스터 폴백) 버그
        return json.dumps({"error": "should-not-be-called"})

    monkeypatch.setattr(server, "_tool_generate_image", _fake_img)

    slides = [{
        "title": "프로젝트 구조",
        "imagePrompt": "isometric folder tree of repository with frontend, backend, tests, docs, scripts, blue icons",
    }]
    out = asyncio.run(server._tool_generate_pptx({"title": "T", "slides": slides}, ""))
    res = json.loads(out)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        stats = _analyze(res["absPath"])
        assert stats["autoshapes"] > 0, "편집 가능한 네이티브 도형이 없음"
        assert stats["pictures"] == 0, f"통짜 래스터가 생성됨(pictures={stats['pictures']})"
        assert called["n"] == 0, "네이티브 변환됐는데 이미지 생성이 호출됨"
    finally:
        _cleanup(res)


def test_pure_photo_prompt_not_converted():
    """P4 — 다이어그램 신호가 없는 순수 사진 프롬프트는 네이티브로 바뀌지 않는다."""
    kind, _ = server._classify_section_diagram(
        "팀 사진", "a photo of a happy team smiling in a sunny office")
    assert kind == "" or kind == "block", f"사진이 구조 다이어그램으로 오분류: {kind}"
    # 'block'이라도 본문이 100자 미만이면 빈 분류여야 한다(원래 동작 보존)
    kind2, _ = server._classify_section_diagram("배경", "sunset over the ocean")
    assert kind2 == "", f"짧은 사진 프롬프트가 다이어그램으로 오분류: {kind2}"


def test_org_chart_layout_and_no_text_truncation(tmp_path):
    """P5 — 루트+평면 자식 트리는 조직도(루트 상단 + 자식 가로)로 그려지고,
    긴 라벨도 잘리지 않는다(자동 축소). 모든 요소는 편집 가능한 도형."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import native_diagram_pptx as nd

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    long_label = "authentication-and-authorization-service"  # 40자, 잘리면 안 됨
    content = f"repository\n  frontend\n  backend\n  {long_label}\n  docs\n  scripts"
    ok = nd.build_native_diagram(slide, "tree", content)
    assert ok is True

    boxes = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    # 루트(1) + 자식(5) = 6 박스
    assert len(boxes) >= 6, f"조직도 박스 부족: {len(boxes)}"
    texts = [sh.text_frame.text for sh in boxes if sh.has_text_frame]
    joined = " ".join(texts)
    assert "repository" in joined, "루트 라벨 없음"
    # 긴 라벨이 잘리지 않고 완전히 보존됨 (이전: [:90]/[:44] 절단)
    assert long_label in joined, f"긴 라벨이 잘림 — texts={texts}"
    # 커넥터(엘보/화살표)도 편집 가능한 도형으로 존재
    connectors = [sh for sh in slide.shapes
                  if "CONNECTOR" in str(sh.shape_type) or sh.shape_type == MSO_SHAPE_TYPE.LINE]
    assert len(connectors) >= 5, f"커넥터 부족: {len(connectors)}"


if __name__ == "__main__":
    raise SystemExit(pytest.main([__file__, "-q"]))
