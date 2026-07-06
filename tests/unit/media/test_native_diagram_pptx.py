"""Regression — 슬라이드 다이어그램이 편집 가능한 네이티브 도형으로 생성된다 (문제2).

근본 요구: 젠스파크/감마처럼, 슬라이드 안의 도표가 통짜 PNG가 아니라 PowerPoint에서
개별 요소(도형/텍스트/화살표)를 수정할 수 있는 네이티브 형식이어야 한다.

수정: native_diagram_pptx.build_native_diagram가 nativeDiagram spec({type, content})을
python-pptx 도형(ROUNDED_RECTANGLE + text_frame + connector)으로 직접 조립.
_tool_generate_pptx가 nativeDiagram을 받으면 add_picture(PNG) 대신 이 빌더를 사용.

Correctness property:
  P1. build_native_diagram(tree/flow/block)은 하나 이상의 도형을 그린다(True 반환).
  P2. nativeDiagram 슬라이드로 생성한 PPTX는 autoshape>0 + Picture==0 (편집 가능, 래스터 없음).
  P3. 각 다이어그램 도형은 편집 가능한 text_frame 텍스트를 가진다.
  P4. nativeDiagram 없는 슬라이드는 기존 동작 보존(도형 다이어그램 미생성).
"""
import os
import sys
import json
import asyncio
import tempfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[3]))
sys.path.insert(0, str(Path(__file__).resolve().parents[3] / "ai_engine"))


def _analyze(pptx_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    prs = Presentation(pptx_path)
    s = {"autoshapes": 0, "with_text": 0, "connectors": 0, "pictures": 0}
    for slide in prs.slides:
        for shp in slide.shapes:
            st = shp.shape_type
            if st == MSO_SHAPE_TYPE.PICTURE:
                s["pictures"] += 1
            elif st == MSO_SHAPE_TYPE.AUTO_SHAPE:
                s["autoshapes"] += 1
                if shp.has_text_frame and shp.text_frame.text.strip():
                    s["with_text"] += 1
            elif "CONNECTOR" in str(st) or st == MSO_SHAPE_TYPE.LINE:
                s["connectors"] += 1
    return s


def _gen(slides, tmp):
    import server
    os.environ["AE_GENERATED_ROOT"] = str(tmp)
    out = asyncio.run(server._tool_generate_pptx(
        {"title": "T", "slides": slides}, ""))
    return json.loads(out)


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


@pytest.mark.unit
@pytest.mark.parametrize("dtype,content,min_shapes", [
    ("tree", "src/\n  a.py\n  b/\n    c.py", 3),
    ("flow", "입력 -> 처리 -> 출력", 3),
    ("block", "프론트\n백엔드\nDB", 3),
])
def test_build_native_diagram_draws_shapes(dtype, content, min_shapes):
    """P1+P3 — 빌더가 편집 가능 도형+텍스트를 그린다."""
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    import native_diagram_pptx as nd

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    ok = nd.build_native_diagram(slide, dtype, content)
    assert ok is True, f"{dtype}: 도형을 그리지 못함"
    shapes = [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]
    assert len(shapes) >= min_shapes, f"{dtype}: 도형 부족 {len(shapes)}"
    # P3 — 텍스트 편집 가능
    texted = [sh for sh in shapes if sh.has_text_frame and sh.text_frame.text.strip()]
    assert len(texted) >= min_shapes, f"{dtype}: 텍스트 도형 부족"


@pytest.mark.unit
def test_build_native_diagram_empty_content_returns_false():
    """빈 콘텐츠 → False (호출자가 PNG 폴백)."""
    from pptx import Presentation
    import native_diagram_pptx as nd
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    assert nd.build_native_diagram(slide, "tree", "") is False
    assert nd.build_native_diagram(slide, "flow", "   ") is False


@pytest.mark.unit
def test_pptx_with_native_diagram_has_no_pictures(tmp_path):
    """P2 — nativeDiagram 슬라이드 PPTX: autoshape>0 + Picture==0 (편집 가능, 래스터 없음)."""
    slides = [
        {"title": "트리", "bullets": ["x"],
         "nativeDiagram": {"type": "tree", "content": "root/\n  a.py\n  b.py"}},
        {"title": "흐름", "bullets": ["y"],
         "nativeDiagram": {"type": "flow", "content": "A -> B -> C"}},
    ]
    res = _gen(slides, tmp_path)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        stats = _analyze(res["absPath"])
        assert stats["autoshapes"] > 0, "편집 가능 도형 없음"
        assert stats["pictures"] == 0, f"통짜 PNG 존재({stats['pictures']}) — 편집 불가 요소"
        assert stats["with_text"] > 0, "텍스트 도형 없음"
    finally:
        _cleanup(res)


@pytest.mark.unit
def test_pptx_without_native_diagram_preserved(tmp_path):
    """P4 — nativeDiagram 없는 슬라이드는 다이어그램 도형 미생성(기존 동작 보존)."""
    slides = [{"title": "텍스트만", "bullets": ["항목1", "항목2"]}]
    res = _gen(slides, tmp_path)
    try:
        assert "error" not in res, f"생성 실패: {res}"
        stats = _analyze(res["absPath"])
        # 다이어그램 도형도 PNG도 없어야 함 (텍스트 placeholder만)
        assert stats["connectors"] == 0
        assert stats["pictures"] == 0
    finally:
        _cleanup(res)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
