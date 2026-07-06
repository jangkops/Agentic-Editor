"""Regression — _execute_tool이 generate_pptx 직접 호출 시 활성 템플릿을 주입한다.

근본 버그(pptx-template-styling 후속): 단일 채팅(run-agent)에서 모델이 generate_pptx를
*직접* 호출하면 _execute_tool이 template_id를 받지 못해 templatePath가 주입되지 않았다.
결과: 빈 Presentation()으로 생성 → 등록된 템플릿(마스터/레이아웃/테마)이 전혀 적용 안 됨.

수정: _execute_tool(template_id) 추가 → generate_pptx 분기에서 _resolve_active_template로
templatePath/templateId/styleProfile을 tool_input에 주입.

Correctness property:
  P1. 유효한 template_id로 generate_pptx 직접 호출 시 응답에 templateId 포함 (적용됨).
  P2. template_id 없이(또는 빈값) 호출하면 응답에 templateId 미포함 (무템플릿 동작 보존).
  P3. 두 경우 모두 디스크에 실재하는 PPTX 파일 생성 (폴백 격리 — 항상 산출).
"""
import os
import sys
import json
import shutil
import tempfile
from pathlib import Path

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[3]))
sys.path.insert(0, str(Path(__file__).resolve().parents[3] / "ai_engine"))


@pytest.fixture
def template_env(tmp_path, monkeypatch):
    """AE_GENERATED_ROOT를 tmp로 설정하고 샘플 템플릿 1개 등록."""
    store = tmp_path / "userdata"
    store.mkdir()
    monkeypatch.setenv("AE_GENERATED_ROOT", str(store))

    import template_manager as tm
    from pptx import Presentation

    # 샘플 .pptx 생성
    sample = tmp_path / "brand.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.save(str(sample))

    store_root = tm.resolve_template_store_root()
    reg = tm.register_template(str(sample), "회귀용 템플릿", store_root)
    assert "error" not in reg, f"템플릿 등록 실패: {reg}"
    return {"templateId": reg["templateId"], "store_root": store_root}


def _gen(template_id):
    import server
    out = server._execute_tool(
        "generate_pptx",
        {"title": "T", "slides": [{"title": "s", "bullets": ["a"]}]},
        project_path="", aws_profile="", bedrock_user="",
        template_id=template_id,
    )
    return json.loads(out)


def _cleanup(res):
    ap = res.get("absPath")
    if ap and os.path.isfile(ap):
        try:
            os.remove(ap)
            if os.path.isfile(ap + ".meta.json"):
                os.remove(ap + ".meta.json")
        except OSError:
            pass


@pytest.mark.unit
def test_execute_tool_injects_template_when_id_present(template_env):
    """P1+P3 — 유효 template_id → 응답에 templateId 포함 + 실재 파일."""
    res = _gen(template_env["templateId"])
    try:
        assert "error" not in res, f"생성 실패: {res}"
        assert res.get("templateId") == template_env["templateId"], \
            "generate_pptx 직접 호출에 템플릿이 주입되지 않음 (빈 Presentation 생성)"
        ap = res.get("absPath")
        assert ap and os.path.isfile(ap) and os.path.getsize(ap) > 0
    finally:
        _cleanup(res)


@pytest.mark.unit
def test_execute_tool_no_template_when_id_absent(template_env):
    """P2+P3 — template_id 없음 → templateId 미포함(무템플릿) + 실재 파일."""
    res = _gen("")
    try:
        assert "error" not in res, f"생성 실패: {res}"
        assert "templateId" not in res, \
            "무템플릿 호출인데 응답에 templateId가 있음 (기존 동작 미보존)"
        ap = res.get("absPath")
        assert ap and os.path.isfile(ap) and os.path.getsize(ap) > 0
    finally:
        _cleanup(res)


@pytest.mark.unit
def test_execute_tool_invalid_template_id_falls_back(template_env):
    """존재하지 않는 template_id → 무템플릿 폴백 (예외 없이 산출, 요구사항 5.4)."""
    res = _gen("nonexistent-template-id-xyz")
    try:
        assert "error" not in res, f"생성 실패(폴백 미작동): {res}"
        # 존재하지 않는 템플릿 → 적용 안 됨 → templateId 미포함
        assert "templateId" not in res
        ap = res.get("absPath")
        assert ap and os.path.isfile(ap) and os.path.getsize(ap) > 0
    finally:
        _cleanup(res)


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
